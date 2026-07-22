# -*- coding: utf-8 -*-
"""生成《下半年单场活动立项计划》汇报 PPT（16:9）· 商务紫主题。
每场一页立项页，便于单场走 OA；不含风险/转化率测算章节。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import plan_data as D

PLUM = RGBColor(0x2E, 0x1F, 0x47)
PURPLE = RGBColor(0x5B, 0x3E, 0x8E)
PURPLE2 = RGBColor(0x76, 0x58, 0xA8)
LILAC = RGBColor(0xC0, 0xAE, 0xE0)
LAV = RGBColor(0xEC, 0xE6, 0xF7)
LAVED = RGBColor(0xDD, 0xD2, 0xEF)
BGT = RGBColor(0xFB, 0xFA, 0xFE)
GOLD = RGBColor(0xC1, 0x9A, 0x3A)
GREY = RGBColor(0x60, 0x5A, 0x6B)
DARK = RGBColor(0x2B, 0x25, 0x36)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TIER_C = {"A": RGBColor(0x86, 0x6F, 0xB4), "B": PURPLE}

FONT = "Microsoft YaHei"
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def _set_font(run, size, bold, color, font=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", font)


def rect(s, x, y, w, h, fill=None, line=None, line_w=None, shadow=False, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w or 1)
    sp.shadow.inherit = False
    if shadow:
        el = sp._element.spPr
        ef = el.makeelement(qn("a:effectLst"), {})
        sh = ef.makeelement(
            qn("a:outerShdw"),
            {"blurRad": "55000", "dist": "25000", "dir": "5400000", "rotWithShape": "0"},
        )
        clr = sh.makeelement(qn("a:srgbClr"), {"val": "7A6BA0"})
        alpha = clr.makeelement(qn("a:alpha"), {"val": "38000"})
        clr.append(alpha)
        sh.append(clr)
        ef.append(sh)
        el.append(ef)
    return sp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0, wrap=True):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for txt, size, bold, color in para:
            r = p.add_run()
            r.text = txt
            _set_font(r, size, bold, color)
    return tb


def bg(s):
    rect(s, 0, 0, SW, SH, fill=BGT)


def header(s, kicker, title, idx=None):
    rect(s, 0, 0, SW, Inches(1.18), fill=PLUM)
    rect(s, 0, Inches(1.18), SW, Pt(3), fill=GOLD)
    rect(s, Inches(0.55), Inches(0.5), Pt(4), Inches(0.5), fill=GOLD)
    text(s, Inches(0.72), Inches(0.18), Inches(11.4), Inches(0.32),
         [[(kicker, 11, True, LILAC)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(0.72), Inches(0.46), Inches(11.4), Inches(0.62),
         [[(title, 24, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    if idx:
        text(s, Inches(11.9), Inches(0.12), Inches(1.1), Inches(0.95),
             [[(idx, 32, True, PURPLE2)]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def footer(s, n):
    rect(s, Inches(0.55), Inches(7.06), Inches(12.25), Pt(0.75), fill=LAVED)
    text(s, Inches(0.55), Inches(7.1), Inches(10), Inches(0.3),
         [[("东方枢纽 × 复旦大学 × 上海市科技企业联合会  |  单场立项计划（8–12 月）", 8, False, GREY)]])
    text(s, Inches(12.2), Inches(7.1), Inches(0.8), Inches(0.3),
         [[(str(n), 8, True, PURPLE)]], align=PP_ALIGN.RIGHT)


# ============================================================ 1. 封面
s = slide()
rect(s, 0, 0, SW, SH, fill=PLUM)
rect(s, 0, 0, SW, Inches(2.55), fill=PURPLE)
rect(s, 0, Inches(2.55), SW, Pt(4), fill=GOLD)
for cx in [Inches(9.55), Inches(10.9), Inches(12.25)]:
    rect(s, cx, Inches(4.55), Inches(1.05), Inches(2.0), fill=None, line=RGBColor(0x53, 0x40, 0x7C), line_w=1.2)
rect(s, Inches(12.25), Inches(4.55), Inches(1.05), Inches(2.0), fill=RGBColor(0x3A, 0x27, 0x5C))
text(s, Inches(0.72), Inches(0.55), Inches(12), Inches(0.5),
     [[("三方战略合作 · 单场活动立项计划", 15, True, LILAC)]])
text(s, Inches(0.72), Inches(1.12), Inches(12), Inches(1.5),
     [[("东方枢纽 A 片区整体办公招商", 36, True, WHITE)],
      [("10 场单场活动 · 独立立项方案", 36, True, WHITE)]], space_after=2)
rect(s, Inches(0.75), Inches(3.02), Inches(2.0), Pt(3), fill=GOLD)
text(s, Inches(0.72), Inches(3.2), Inches(12), Inches(0.5),
     [[("复旦大学  ·  上海市科技企业联合会  ·  东方枢纽", 15, True, WHITE)]])
text(s, Inches(0.72), Inches(3.78), Inches(11.8), Inches(0.5),
     [[("8–12 月每月 2 场  |  单场独立走 OA  |  围绕六大产业  |  市级 + 浦东新区资源联动", 12.5, False, LILAC)]])
text(s, Inches(0.72), Inches(6.62), Inches(12), Inches(0.5),
     [[("2026 年 8–12 月  |  内部审阅 & 客户呈送 / 单场立项通用版", 11, False, LILAC)]])

# ============================================================ 2. 目录
s = slide()
bg(s)
header(s, "AGENDA", "汇报目录")
items = [
    ("01", "合作背景与三方定位", "谁在合作 · 各自角色"),
    ("02", "政府资源背书矩阵", "市级 + 浦东新区联动"),
    ("03", "招商标的 · A 片区办公项目", "四大产品线 · 销售/租赁"),
    ("04", "六大产业与月度排期", "8–12 月 · 每月 2 场"),
    ("05", "10 场活动总览", "立项编号 · 主题 · 报价"),
    ("06", "单场立项详情（10 页）", "每场单独成案 · 便于 OA"),
    ("07", "报价体系", "两档模型 · 全年预算"),
    ("08", "价值要点与执行下一步", "六大产业利好 · 单场 OA"),
]
cx = [Inches(0.6), Inches(6.9)]
for i, (num, t, sub) in enumerate(items):
    col, row = i % 2, i // 2
    x, y = cx[col], Inches(1.5) + row * Inches(1.3)
    w, h = Inches(5.85), Inches(1.1)
    rect(s, x, y, w, h, fill=WHITE, line=LAVED, line_w=1, shadow=True)
    rect(s, x, y, Inches(1.0), h, fill=PURPLE)
    text(s, x, y, Inches(1.0), h, [[(num, 22, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + Inches(1.18), y + Inches(0.22), w - Inches(1.3), Inches(0.4), [[(t, 14, True, PLUM)]])
    text(s, x + Inches(1.18), y + Inches(0.62), w - Inches(1.3), Inches(0.35), [[(sub, 10, False, GREY)]])
footer(s, 2)

# ============================================================ 3. 三方定位
s = slide()
bg(s)
header(s, "背景 · BACKGROUND", "合作背景与三方定位", "01")
text(s, Inches(0.72), Inches(1.38), Inches(12), Inches(0.45),
     [[("以“上海市级 + 浦东新区”政府资源联动为核心；每场活动单独立项、单独报价、单独走 OA。", 12, False, GREY)]])
colors = [PURPLE, PURPLE2, PLUM]
xs = [Inches(0.6), Inches(4.87), Inches(9.14)]
for i, (name, role, desc) in enumerate(D.PARTIES):
    x, y, w, h = xs[i], Inches(2.0), Inches(3.87), Inches(3.85)
    rect(s, x, y, w, h, fill=WHITE, line=LAVED, line_w=1, shadow=True)
    rect(s, x, y, w, Inches(1.05), fill=colors[i])
    text(s, x + Inches(0.22), y + Inches(0.14), w - Inches(0.44), Inches(0.5), [[(name, 15, True, WHITE)]])
    text(s, x + Inches(0.22), y + Inches(0.62), w - Inches(0.44), Inches(0.35), [[(role, 10.5, True, LILAC)]])
    text(s, x + Inches(0.24), y + Inches(1.24), w - Inches(0.48), h - Inches(1.4),
         [[(desc, 11.5, False, DARK)]], line_spacing=1.28)
rect(s, Inches(0.6), Inches(6.05), Pt(4), Inches(0.55), fill=GOLD)
text(s, Inches(0.78), Inches(6.05), Inches(12), Inches(0.55),
     [[("备选/补充科技组织：", 10.5, True, PURPLE), ("、".join(D.ALT_TECH_ORGS), 10.5, False, GREY)]],
     anchor=MSO_ANCHOR.MIDDLE)
footer(s, 3)

# ============================================================ 4. 政府资源
s = slide()
bg(s)
header(s, "资源 · GOVERNMENT", "政府资源背书矩阵（市级 + 浦东新区）", "02")
text(s, Inches(0.72), Inches(1.4), Inches(12), Inches(0.4),
     [[("核心逻辑：市级资源与浦东新区资源相结合，商务委与投资促进部门共同参与背书。", 12, True, PLUM)]])
groups = list(D.GOV_RESOURCES.items())
gcolors = [PURPLE, PLUM]
for i, (group, gitems) in enumerate(groups):
    x, y, w, h = Inches(0.6) + i * Inches(6.35), Inches(2.05), Inches(6.0), Inches(3.35)
    rect(s, x, y, w, h, fill=WHITE, line=LAVED, line_w=1, shadow=True)
    rect(s, x, y, w, Inches(0.72), fill=gcolors[i])
    text(s, x + Inches(0.25), y, w - Inches(0.4), Inches(0.72), [[(group, 16, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    for j, it in enumerate(gitems):
        yy = y + Inches(0.95) + j * Inches(0.72)
        rect(s, x + Inches(0.25), yy, w - Inches(0.5), Inches(0.58), fill=LAV)
        rect(s, x + Inches(0.25), yy, Pt(4), Inches(0.58), fill=GOLD)
        text(s, x + Inches(0.45), yy, w - Inches(0.7), Inches(0.58), [[(it, 12, True, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(6.35), Inches(3.3), Inches(0.65), Inches(0.7), [[("＋", 30, True, GOLD)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
y = Inches(5.7)
rect(s, Inches(0.6), y, Inches(12.35), Inches(0.95), fill=PLUM)
rect(s, Inches(0.6), y, Pt(5), Inches(0.95), fill=GOLD)
text(s, Inches(0.9), y, Inches(11.8), Inches(0.95), [[(D.GOV_TAGLINE, 13, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 4)

# ============================================================ 5. 招商标的
s = slide()
bg(s)
header(s, "标的 · PROJECT", "招商标的 · A 片区整体办公项目", "03")
P = D.PROJECT
text(s, Inches(0.72), Inches(1.35), Inches(12.3), Inches(0.35), [[(P["position"], 11, False, GREY)]])
rect(s, Inches(0.6), Inches(1.78), Inches(12.35), Inches(0.5), fill=LAV)
rect(s, Inches(0.6), Inches(1.78), Pt(5), Inches(0.5), fill=GOLD)
text(s, Inches(0.82), Inches(1.78), Inches(12), Inches(0.5),
     [[("体量：", 11, True, PLUM), (P["area"] + "  ", 11, True, GOLD),
       ("（招商标的为 A 片区办公项目，非“133 万方”）", 10, False, GREY)]],
     anchor=MSO_ANCHOR.MIDDLE)
for i, (n, d) in enumerate(P["product_lines"]):
    col, row = i % 2, i // 2
    x = Inches(0.6) + col * Inches(6.3)
    y = Inches(2.45) + row * Inches(1.15)
    w, h = Inches(6.0), Inches(1.0)
    rect(s, x, y, w, h, fill=WHITE, line=LAVED, line_w=1, shadow=True)
    rect(s, x, y, Inches(0.14), h, fill=[PURPLE, PURPLE2, PLUM, TIER_C["A"]][i])
    text(s, x + Inches(0.3), y + Inches(0.08), w - Inches(0.4), Inches(0.32),
         [[(f"产品线 {i+1} · {n}", 12, True, PURPLE)]])
    text(s, x + Inches(0.3), y + Inches(0.42), w - Inches(0.5), Inches(0.55),
         [[(d, 9.2, False, DARK)]], line_spacing=1.05)
y = Inches(4.9)
rect(s, Inches(0.6), y, Inches(12.35), Inches(1.7), fill=WHITE, line=LAVED, line_w=1, shadow=True)
rect(s, Inches(0.6), y, Inches(0.14), Inches(1.7), fill=GOLD)
text(s, Inches(0.85), y + Inches(0.1), Inches(12), Inches(0.3), [[("销售 / 租赁模式", 12, True, PLUM)]])
runs = [[("· ", 9.5, True, GOLD), (x, 9.5, False, DARK)] for x in P["model"]]
text(s, Inches(0.87), y + Inches(0.42), Inches(12), Inches(0.75), runs, space_after=2, line_spacing=1.0)
text(s, Inches(0.85), y + Inches(1.25), Inches(12), Inches(0.35),
     [[("活动↔产品匹配：", 9.5, True, PURPLE), (P["match"], 9.5, False, GREY)]])
footer(s, 5)

# ============================================================ 6. 六大产业 + 月度排期
s = slide()
bg(s)
header(s, "产业 · 排期", "六大核心产业与月度排期（8–12 月）", "04")
icons = ["高服", "国贸", "民航", "医药", "工联", "绿能"]
for i, (name, desc) in enumerate(D.INDUSTRIES):
    col = i % 3
    row = i // 3
    x = Inches(0.55) + col * Inches(2.85)
    y = Inches(1.4) + row * Inches(1.55)
    w, h = Inches(2.7), Inches(1.4)
    rect(s, x, y, w, h, fill=WHITE, line=LAVED, line_w=1, shadow=True)
    rect(s, x, y, w, Inches(0.1), fill=PURPLE)
    text(s, x + Inches(0.12), y + Inches(0.2), w - Inches(0.2), Inches(0.35),
         [[(f"{icons[i]} · {name}", 11, True, PLUM)]])
    text(s, x + Inches(0.12), y + Inches(0.6), w - Inches(0.24), Inches(0.7),
         [[(desc, 9, False, GREY)]], line_spacing=1.05)
# 右侧月度
rect(s, Inches(9.2), Inches(1.4), Inches(3.6), Inches(5.2), fill=WHITE, line=LAVED, line_w=1, shadow=True)
rect(s, Inches(9.2), Inches(1.4), Inches(3.6), Inches(0.55), fill=PLUM)
text(s, Inches(9.2), Inches(1.4), Inches(3.6), Inches(0.55),
     [[("月度排期 · 每月 2 场", 13, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
for i, (month, cnt, note) in enumerate(D.MONTH_PLAN):
    y = Inches(2.1) + i * Inches(0.85)
    text(s, Inches(9.4), y, Inches(3.2), Inches(0.3), [[(f"{month}  ·  {cnt} 场", 12, True, PURPLE)]])
    text(s, Inches(9.4), y + Inches(0.3), Inches(3.2), Inches(0.45), [[(note, 9, False, GREY)]])
footer(s, 6)

# ============================================================ 7. 10场总览
s = slide()
bg(s)
header(s, "总览 · OVERVIEW", "10 场单场活动总览（独立立项）", "05")
cols = [("#", 0.4), ("立项编号", 1.35), ("时间", 1.7), ("活动主题", 4.0), ("产业板块", 1.9), ("规模", 1.2), ("档", 0.45), ("报价", 0.7)]
x0, y0, rowh = Inches(0.5), Inches(1.35), Inches(0.48)
cxs, acc = [], x0
for name, w in cols:
    cxs.append((acc, Inches(w)))
    acc += Inches(w)
for (cx, cw), (name, _) in zip(cxs, cols):
    rect(s, cx, y0, cw, rowh, fill=PLUM)
    text(s, cx, y0, cw, rowh, [[(name, 9.5, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
for i, a in enumerate(D.ACTIVITIES):
    y = y0 + rowh * (i + 1)
    fill = LAV if i % 2 else WHITE
    tc = a["tier"].split(" ")[0]
    vals = [str(a["no"]), a["oa_code"], a["date"], a["title"], a["sector"],
            a["scale"].replace(" · ", "·"), tc, str(a["price"])]
    for j, ((cx, cw), v) in enumerate(zip(cxs, vals)):
        rect(s, cx, y, cw, rowh, fill=fill, line=LAVED, line_w=0.5)
        al = PP_ALIGN.LEFT if j in (3, 4) else PP_ALIGN.CENTER
        col = TIER_C[tc] if j == 6 else (GOLD if j == 7 else DARK)
        pad = Inches(0.06) if j in (3, 4) else 0
        text(s, cx + pad, y, cw - pad, rowh, [[(v, 8.5, j in (6, 7), col)]],
             align=al, anchor=MSO_ANCHOR.MIDDLE)
yt = y0 + rowh * (len(D.ACTIVITIES) + 1)
rect(s, cxs[0][0], yt, cxs[6][0] + cxs[6][1] - cxs[0][0], rowh, fill=LILAC)
text(s, cxs[0][0], yt, cxs[6][0] + cxs[6][1] - cxs[0][0], rowh,
     [[("10 场合计（单场立项累加）", 10, True, PLUM)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rect(s, cxs[7][0], yt, cxs[7][1], rowh, fill=GOLD)
text(s, cxs[7][0], yt, cxs[7][1], rowh, [[(str(D.TOTAL_PRICE), 11, True, WHITE)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.3),
     [[("说明：7 月不排场；最早一场 8 月上旬。每场单独编号（DFSN-2026-XX），可单独提交 OA。", 8.5, False, GREY)]])

# ============================================================ 8–17. 单场立项页（每场一页）
def oa_slide(a, page_no):
    s = slide()
    bg(s)
    tc = a["tier"].split(" ")[0]
    accent = TIER_C[tc]
    header(s, f"单场立项 · {a['oa_code']}", f"NO.{a['no']}  {a['title']}", "06")
    # 左栏：基本信息
    rect(s, Inches(0.5), Inches(1.4), Inches(4.3), Inches(5.2), fill=WHITE, line=LAVED, line_w=1, shadow=True)
    rect(s, Inches(0.5), Inches(1.4), Inches(4.3), Inches(0.5), fill=accent)
    text(s, Inches(0.5), Inches(1.4), Inches(4.3), Inches(0.5),
         [[("基本信息 · OA 要件", 13, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    meta = [
        ("立项编号", a["oa_code"]),
        ("拟定时间", a["date"]),
        ("所属月份", a["month"]),
        ("产业板块", a["sector"]),
        ("形式规模", a["scale"]),
        ("场地建议", a["venue"]),
        ("报价档位", a["tier"].split(" · ")[0]),
    ]
    for i, (k, v) in enumerate(meta):
        y = Inches(2.05) + i * Inches(0.42)
        text(s, Inches(0.7), y, Inches(1.3), Inches(0.38), [[(k, 10, True, GREY)]], anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(2.0), y, Inches(2.6), Inches(0.38), [[(v, 10.5, True, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
    # 报价大字
    rect(s, Inches(0.7), Inches(5.1), Inches(3.9), Inches(1.2), fill=LAV)
    text(s, Inches(0.7), Inches(5.2), Inches(3.9), Inches(0.35),
         [[("单场报价（万元）", 11, False, GREY)]], align=PP_ALIGN.CENTER)
    text(s, Inches(0.7), Inches(5.55), Inches(3.9), Inches(0.6),
         [[(str(a["price"]), 28, True, accent)]], align=PP_ALIGN.CENTER)

    # 右栏：内容
    x, w = Inches(5.05), Inches(7.7)
    rect(s, x, Inches(1.4), w, Inches(5.2), fill=WHITE, line=LAVED, line_w=1, shadow=True)
    text(s, x + Inches(0.25), Inches(1.55), w - Inches(0.4), Inches(0.3),
         [[("■ 内容建议", 12, True, accent)]])
    runs = [[("· ", 10, True, GOLD), (ln, 10, False, DARK)] for ln in a["content"]]
    text(s, x + Inches(0.3), Inches(1.9), w - Inches(0.5), Inches(1.2), runs, space_after=3, line_spacing=1.05)

    text(s, x + Inches(0.25), Inches(3.2), w - Inches(0.4), Inches(0.3),
         [[("■ 拟邀嘉宾 / 资源", 12, True, accent)]])
    text(s, x + Inches(0.3), Inches(3.55), w - Inches(0.5), Inches(0.7),
         [[("、".join(a["guests"]), 10.5, False, DARK)]], line_spacing=1.1)

    text(s, x + Inches(0.25), Inches(4.35), w - Inches(0.4), Inches(0.3),
         [[("■ 招商衔接价值", 12, True, accent)]])
    runs = [[("· ", 10, True, GOLD), (ln, 10, False, DARK)] for ln in a["invest"]]
    text(s, x + Inches(0.3), Inches(4.7), w - Inches(0.5), Inches(0.85), runs, space_after=2, line_spacing=1.05)

    text(s, x + Inches(0.25), Inches(5.7), w - Inches(0.4), Inches(0.55),
         [[("立项说明：", 10, True, PURPLE), (a["oa_note"], 10, False, GREY)]], line_spacing=1.05)
    footer(s, page_no)


pg = 8
for a in D.ACTIVITIES:
    oa_slide(a, pg)
    pg += 1

# ============================================================ 报价体系
s = slide()
bg(s)
header(s, "报价 · QUOTE", "报价体系 · 两档模型与全年预算", "07")
tier_keys = list(D.TIERS.keys())
for i, k in enumerate(tier_keys):
    x = Inches(0.7) + i * Inches(6.2)
    y, w, h = Inches(1.4), Inches(5.85), Inches(3.4)
    rect(s, x, y, w, h, fill=WHITE, line=LAVED, line_w=1, shadow=True)
    rect(s, x, y, w, Inches(0.9), fill=TIER_C[k.split(" ")[0]])
    parts = k.split(" · ")
    text(s, x + Inches(0.2), y + Inches(0.1), w - Inches(0.3), Inches(0.4), [[(parts[0], 15, True, WHITE)]])
    text(s, x + Inches(0.2), y + Inches(0.5), w - Inches(0.3), Inches(0.35),
         [[(parts[1] if len(parts) > 1 else "", 11, False, LILAC)]])
    text(s, x + Inches(0.2), y + Inches(1.05), w - Inches(0.3), Inches(0.5),
         [[(f"{D.TIER_TOTAL[k]} ", 26, True, PURPLE), ("万元/场", 12, True, GREY)]])
    runs = [[(item, 10, False, DARK), (f"  {v}", 10, True, GOLD)]
            for item, v in zip(D.COST_ITEMS, D.TIERS[k])]
    text(s, x + Inches(0.25), y + Inches(1.65), w - Inches(0.4), Inches(1.6), runs, space_after=2)

tier_counts = {}
for a in D.ACTIVITIES:
    tier_counts[a["tier"]] = tier_counts.get(a["tier"], 0) + 1
y = Inches(5.05)
rect(s, Inches(0.5), y, Inches(12.33), Inches(1.6), fill=LAV, line=LAVED, line_w=0.75)
text(s, Inches(0.7), y + Inches(0.12), Inches(12), Inches(0.35),
     [[("全年预算汇总（单场立项累加）", 13, True, PLUM),
       ("    本期 10 场均为 A/B 档；每场可单独审批、单独结算", 10, False, GREY)]])
boxes = []
for k in tier_keys:
    cnt = tier_counts.get(k, 0)
    boxes.append((f"{k.split(' ')[0]} 档 × {cnt} 场", round(cnt * D.TIER_TOTAL[k], 1),
                  TIER_C[k.split(" ")[0]], False))
boxes.append(("10 场合计", D.TOTAL_PRICE, PLUM, True))
n = len(boxes)
gap = Inches(0.25)
bw = (Inches(11.9) - gap * (n - 1)) / n
for i, (lab, val, colr, is_total) in enumerate(boxes):
    x = Inches(0.7) + i * (bw + gap)
    yy = y + Inches(0.55)
    rect(s, x, yy, bw, Inches(0.85), fill=(PLUM if is_total else WHITE),
         line=(None if is_total else colr), line_w=1.2)
    text(s, x, yy + Inches(0.08), bw, Inches(0.35),
         [[(lab, 11, True, (LILAC if is_total else colr))]], align=PP_ALIGN.CENTER)
    text(s, x, yy + Inches(0.42), bw, Inches(0.4),
         [[(f"{val} 万元", 14, True, (WHITE if is_total else PLUM))]], align=PP_ALIGN.CENTER)
footer(s, pg)
pg += 1

# ============================================================ 价值要点（定性，无转化率/风险）
s = slide()
bg(s)
header(s, "价值 · VALUE", "招商价值要点（围绕六大产业）", "08")
text(s, Inches(0.72), Inches(1.35), Inches(12), Inches(0.4),
     [[("10 场单场活动，全面覆盖东方枢纽既定六大产业板块，服务 A 片区整体办公招商。", 12, True, PLUM)]])
for i, (name, desc) in enumerate(D.VALUE_PILLARS):
    col, row = i % 3, i // 3
    x = Inches(0.55) + col * Inches(4.15)
    y = Inches(1.95) + row * Inches(2.2)
    w, h = Inches(3.95), Inches(2.0)
    rect(s, x, y, w, h, fill=WHITE, line=LAVED, line_w=1, shadow=True)
    rect(s, x, y, w, Inches(0.55), fill=PURPLE)
    text(s, x + Inches(0.2), y, w - Inches(0.3), Inches(0.55),
         [[(f"{i+1}. {name}", 13, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + Inches(0.22), y + Inches(0.75), w - Inches(0.4), Inches(1.1),
         [[(desc, 11.5, False, DARK)]], line_spacing=1.2)
footer(s, pg)
pg += 1

# ============================================================ 执行 + 下一步
s = slide()
bg(s)
header(s, "执行 · NEXT", "单场 OA 执行机制与下一步", "08")
for i, (name, desc) in enumerate(D.EXECUTION):
    col, row = i % 3, i // 3
    x = Inches(0.55) + col * Inches(4.15)
    y = Inches(1.4) + row * Inches(1.55)
    w, h = Inches(3.95), Inches(1.4)
    rect(s, x, y, w, h, fill=WHITE, line=LAVED, line_w=0.75, shadow=True)
    rect(s, x, y, Inches(0.7), h, fill=PURPLE)
    text(s, x, y, Inches(0.7), h, [[(f"{i+1:02d}", 16, True, WHITE)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + Inches(0.85), y + Inches(0.15), w - Inches(1.0), Inches(0.35), [[(name, 12, True, PLUM)]])
    text(s, x + Inches(0.85), y + Inches(0.55), w - Inches(1.05), Inches(0.75),
         [[(desc, 9.5, False, DARK)]], line_spacing=1.1)

# 下一步条
y = Inches(4.7)
rect(s, Inches(0.55), y, Inches(12.25), Inches(1.95), fill=PLUM)
text(s, Inches(0.8), y + Inches(0.15), Inches(11.8), Inches(0.4),
     [[("下一步 · 建议从第 1 场启动立项", 14, True, GOLD)]])
steps = [
    "① 确认 8 月上旬首场（进境关外政策闭门会）主题、时间与 30 人邀约名单",
    "② 按单场立项表提交 OA 要件：名单 / 人数 / 背景 / 预算 / 时间",
    "③ 经指定策划供应商签约支付，落实场地搭建通行证",
    "④ 8–12 月每月 2 场滚动执行，一场一立、一场一结",
]
for i, t in enumerate(steps):
    text(s, Inches(0.85), y + Inches(0.55) + i * Inches(0.32), Inches(11.7), Inches(0.3),
         [[(t, 11, False, WHITE)]])
footer(s, pg)

out = "东方枢纽三方合作计划_汇报PPT.pptx"
prs.save(out)
print(f"PPT 已生成：{out}  (页数: {len(prs.slides._sldIdLst)})")
