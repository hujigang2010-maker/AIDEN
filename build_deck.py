# -*- coding: utf-8 -*-
"""
生成《四方联动 · 共建"AI+云计算"产业生态共同体》合作方案 PPT。
运行: python3 build_deck.py  ->  输出 合作方案_AI云计算产业生态共同体.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------
# 主题配色
# ----------------------------------------------------------------------------
NAVY   = RGBColor(0x0B, 0x1F, 0x3A)   # 深海军蓝（主背景）
BLUE   = RGBColor(0x1B, 0x4F, 0x9C)   # 主蓝
CYAN   = RGBColor(0x16, 0xB8, 0xC6)   # 科技青（强调）
GOLD   = RGBColor(0xF2, 0xB6, 0x32)   # 金（点睛）
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xE8, 0xEE, 0xF6)   # 浅灰蓝
GREY   = RGBColor(0x5B, 0x6B, 0x82)
DARKTX = RGBColor(0x1A, 0x26, 0x38)

FONT   = "微软雅黑"

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# ----------------------------------------------------------------------------
# 工具函数
# ----------------------------------------------------------------------------
def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, x, y, w, h, color, line=False, line_color=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line and line_color is not None:
        sp.line.color.rgb = line_color
        sp.line.width = Pt(1)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def textbox(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            space_after=Pt(6), line_spacing=1.05, wrap=True):
    """runs: list of paragraphs; each paragraph is list of (text, size, color, bold)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = space_after
        p.line_spacing = line_spacing
        for (text, size, color, bold) in para:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = FONT
            rPr = r.font._rPr
            ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT})
            rPr.append(ea)
    return tb


def section_header(slide, kicker, title, idx=None):
    """顶部标题区（用于内容页）"""
    set_bg(slide, WHITE)
    rect(slide, 0, 0, EMU_W, Inches(1.25), NAVY)
    rect(slide, 0, Inches(1.25), EMU_W, Inches(0.06), CYAN)
    # 左侧色条
    rect(slide, Inches(0.0), 0, Inches(0.18), Inches(1.25), GOLD)
    runs = []
    if kicker:
        runs.append([(kicker, 12, CYAN, True)])
    runs.append([(title, 26, WHITE, True)])
    textbox(slide, Inches(0.55), Inches(0.16), Inches(11.5), Inches(1.0), runs,
            anchor=MSO_ANCHOR.MIDDLE, space_after=Pt(2))
    if idx:
        textbox(slide, Inches(11.8), Inches(0.16), Inches(1.3), Inches(1.0),
                [[(idx, 40, RGBColor(0x2C, 0x44, 0x6E), True)]],
                align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def bullet(slide, x, y, w, h, items, gap=Pt(10), size=15, lead_color=CYAN):
    """items: list of (lead, body) ; lead 加粗着色"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (lead, body) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = gap
        p.line_spacing = 1.12
        r0 = p.add_run(); r0.text = "▍"; r0.font.size = Pt(size); r0.font.color.rgb = lead_color; r0.font.bold = True; r0.font.name = FONT
        if lead:
            r1 = p.add_run(); r1.text = lead + "  "; r1.font.size = Pt(size); r1.font.color.rgb = DARKTX; r1.font.bold = True; r1.font.name = FONT
        r2 = p.add_run(); r2.text = body; r2.font.size = Pt(size); r2.font.color.rgb = GREY; r2.font.bold = False; r2.font.name = FONT
    return tb


def card(slide, x, y, w, h, title, lines, accent=BLUE, title_color=WHITE):
    rect(slide, x, y, w, h, LIGHT)
    rect(slide, x, y, w, Inches(0.62), accent)
    textbox(slide, x + Inches(0.18), y, w - Inches(0.3), Inches(0.62),
            [[(title, 14, title_color, True)]], anchor=MSO_ANCHOR.MIDDLE)
    body = [[(l, 12.5, DARKTX, False)] for l in lines]
    textbox(slide, x + Inches(0.18), y + Inches(0.72), w - Inches(0.34), h - Inches(0.8),
            body, space_after=Pt(5), line_spacing=1.1)


def style_table(table, header_fill=BLUE, header_color=WHITE, body_size=12, header_size=13):
    tbl = table._tbl
    # remove default style banding by setting first row only
    for r_idx, row in enumerate(table.rows):
        for c_idx, cell in enumerate(row.cells):
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.10)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            tf = cell.text_frame
            for p in tf.paragraphs:
                for run in p.runs:
                    run.font.name = FONT
                    if r_idx == 0:
                        run.font.size = Pt(header_size); run.font.bold = True; run.font.color.rgb = header_color
                    else:
                        run.font.size = Pt(body_size); run.font.color.rgb = DARKTX
            if r_idx == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r_idx % 2 else RGBColor(0xED, 0xF2, 0xF9)


def fill_table(table, data):
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            table.cell(r, c).text = val


# ============================================================================
# 幻灯片 1 — 封面
# ============================================================================
s = add_slide()
set_bg(s, NAVY)
rect(s, 0, 0, Inches(0.22), EMU_H, CYAN)
rect(s, Inches(0.22), 0, Inches(0.08), EMU_H, GOLD)
# 装饰圆
deco = rect(s, Inches(9.3), Inches(-1.4), Inches(5.2), Inches(5.2), BLUE, shape=MSO_SHAPE.OVAL)
deco.fill.fore_color.rgb = RGBColor(0x12, 0x33, 0x5E)
deco2 = rect(s, Inches(10.6), Inches(3.0), Inches(4.6), Inches(4.6), BLUE, shape=MSO_SHAPE.OVAL)
deco2.fill.fore_color.rgb = RGBColor(0x0F, 0x2A, 0x4D)

textbox(s, Inches(0.7), Inches(1.0), Inches(11), Inches(0.5),
        [[("战略合作方案 · 内部沟通版", 14, CYAN, True)]])
textbox(s, Inches(0.7), Inches(1.55), Inches(11.6), Inches(2.4),
        [[("五方共建", 48, WHITE, True)],
         [("“AI + 云计算” 产业生态共同体", 38, WHITE, True)]],
        space_after=Pt(8), line_spacing=1.02)
rect(s, Inches(0.75), Inches(3.8), Inches(3.6), Inches(0.05), GOLD)
textbox(s, Inches(0.7), Inches(3.98), Inches(11.9), Inches(1.5),
        [[("面向 · 上海市云计算创新基地（大学路 · 五角场创新创业学院）", 15, LIGHT, False)],
         [("两大诉求：轻资产拓展 × 招商对接", 15, GOLD, True)],
         [("新引擎：腾讯云算力底座 + OPC 超级孵化闭环 → 产业闭环", 15, CYAN, True)]],
        space_after=Pt(5))
textbox(s, Inches(0.7), Inches(6.45), Inches(12.2), Inches(0.8),
        [[("复旦大学住房政策研究中心 · 人工智能商业化落地组委会 · 上海市杨浦区科技企业联合会 · 腾讯云", 12, GREY, False)]])


# ============================================================================
# 幻灯片 2 — 一句话定调
# ============================================================================
s = add_slide()
set_bg(s, NAVY)
rect(s, 0, Inches(2.4), EMU_W, Inches(0.06), CYAN)
textbox(s, Inches(1.0), Inches(0.9), Inches(11.3), Inches(0.6),
        [[("先把这件事结构化", 16, CYAN, True)]])
textbox(s, Inches(1.0), Inches(1.5), Inches(11.3), Inches(1.0),
        [[("这次不是简单谈一单合作 —", 30, WHITE, True)]])
textbox(s, Inches(1.0), Inches(2.7), Inches(11.3), Inches(2.4),
        [[("而是用 “资源整合 + 资产运营 + 产业导入”", 32, GOLD, True)],
         [("共建一个可复制的 AI 产业空间运营新模式。", 32, WHITE, True)]],
        space_after=Pt(10), line_spacing=1.15)
textbox(s, Inches(1.0), Inches(5.35), Inches(11.5), Inches(1.6),
        [[("对方提供“资源平台”（空间 + 政府背书 + 双创资源），腾讯云提供“算力 + 技术底座”；", 16, LIGHT, False)],
         [("我们补上“内容 + 运营 + 产业导入 + 商业化变现”，并用「超级孵化器」批量孵化 OPC，", 16, LIGHT, False)],
         [("（AI 一人公司 / 超级个体）—— 让所有 OPC 在生态链里完成闭环。", 16, GOLD, True)]],
        space_after=Pt(6))


# ============================================================================
# 幻灯片 3 — 资源禀赋（双方）
# ============================================================================
s = add_slide()
section_header(s, "RESOURCES · 资源禀赋", "五方资源盘点：空间 × 内容 × 算力，三位一体", "01")
# 左：对方
rect(s, Inches(0.55), Inches(1.5), Inches(5.9), Inches(4.35), LIGHT)
rect(s, Inches(0.55), Inches(1.5), Inches(5.9), Inches(0.55), BLUE)
textbox(s, Inches(0.75), Inches(1.5), Inches(5.6), Inches(0.55),
        [[("对方 · 资源平台（壳）", 15, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
bullet(s, Inches(0.75), Inches(2.2), Inches(5.5), Inches(3.6), [
    ("云计算创新基地：", "产业承载空间 + 企业集聚载体，孵化器/加速器/总部基地/人才实训“四位一体”。"),
    ("五角场创新创业学院：", "双创资源连接器，“募课”模式、长三角科创板预备营等全周期服务。"),
    ("缺口：", "产业内容不足 + 运营能力缺失，物业易空置。"),
], size=13, lead_color=BLUE)
# 右：我方
rect(s, Inches(6.85), Inches(1.5), Inches(5.9), Inches(4.35), RGBColor(0xEA, 0xF6, 0xF8))
rect(s, Inches(6.85), Inches(1.5), Inches(5.9), Inches(0.55), CYAN)
textbox(s, Inches(7.05), Inches(1.5), Inches(5.6), Inches(0.55),
        [[("我方 · 内容运营方（芯）", 15, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
bullet(s, Inches(7.05), Inches(2.2), Inches(5.5), Inches(3.6), [
    ("复旦住房政策研究中心：", "城市更新、存量盘活、产业社区规划与政策背书。"),
    ("AI 商业化落地组委会：", "5·22 峰会汇聚 AI 全产业链企业与投融资资源。"),
    ("杨浦区科技企业联合会：", "覆盖七大战新产业会员企业网络，直接招商渠道。"),
    ("您本人：", "跨组织协同纽带 + 轻资产拓展操盘能力（关键）。"),
], size=13, lead_color=CYAN)
# 底部：腾讯云
rect(s, Inches(0.55), Inches(6.0), Inches(12.2), Inches(0.95), NAVY)
rect(s, Inches(0.55), Inches(6.0), Inches(0.16), Inches(0.95), GOLD)
textbox(s, Inches(0.95), Inches(6.05), Inches(11.6), Inches(0.85),
        [[("新伙伴 · 腾讯云（算力与技术底座）：", 14, GOLD, True),
          ("算力补贴 + 混元大模型/API + 云开发技术中台 + 初创扶持计划 —— 补齐“算力+技术”底座，支撑 OPC 规模化孵化。", 14, WHITE, False)]],
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)


# ============================================================================
# 幻灯片 4 — 诉求 → 本质问题
# ============================================================================
s = add_slide()
section_header(s, "INSIGHT · 核心判断", "两大诉求的背后，是“内容 + 运营”的缺口", "02")
# 两个诉求卡
card(s, Inches(0.55), Inches(1.7), Inches(5.9), Inches(2.2), "诉求 1 · 轻资产拓展", [
    "以大学路五角场创新创业学院为基础，", "寻找高空置率物业进行合作。",
    "痛点：业主方凭什么把盘子交出来？"], accent=BLUE)
card(s, Inches(6.85), Inches(1.7), Inches(5.9), Inches(2.2), "诉求 2 · 招商对接", [
    "为产业办公用地寻找招商租客，", "吸引企业入驻。",
    "痛点：纯中介式招商无护城河、匹配度低。"], accent=BLUE)
# 箭头/结论
rect(s, Inches(0.55), Inches(4.2), Inches(12.2), Inches(2.5), NAVY)
rect(s, Inches(0.55), Inches(4.2), Inches(0.16), Inches(2.5), GOLD)
textbox(s, Inches(1.0), Inches(4.45), Inches(11.5), Inches(0.6),
        [[("本质问题不是物业和招商，而是 ——", 16, CYAN, True)]])
textbox(s, Inches(1.0), Inches(5.05), Inches(11.5), Inches(1.6),
        [[("产业内容不足  +  运营能力缺失。", 26, WHITE, True)],
         [("我们恰好能补这块能力，并把“找物业 / 找租客”升级为“资产重塑 + 生态导入”。", 16, LIGHT, False)]],
        space_after=Pt(10))


# ============================================================================
# 幻灯片 5 — 合作点1 轻资产拓展
# ============================================================================
s = add_slide()
section_header(s, "合作点 1 · 轻资产拓展", "从“找空置”升级为“资产重塑与联营”", "03")
card(s, Inches(0.55), Inches(1.65), Inches(3.9), Inches(2.55), "① 政产学研赋能包", [
    "复旦中心出具区域产业焕新 / 城市更新评估，",
    "为国企、政府平台业主方提供理论背书，",
    "大幅提升交盘信任度与合作意愿。"], accent=BLUE)
card(s, Inches(4.7), Inches(1.65), Inches(3.9), Inches(2.55), "② 复用操盘 · 品牌+运营双输出", [
    "复用杨浦科技商业广场出租率",
    "70% → 90%+ 的盘活逻辑；",
    "对方输出品牌+产业资源，",
    "我方输出资产调改与运营体系。"], accent=BLUE)
card(s, Inches(8.85), Inches(1.65), Inches(3.9), Inches(2.55), "③ 主题先导区 · 产品化", [
    "不做传统二房东，",
    "包装为“AI+产业加速器空间”，",
    "如挂牌“杨浦大模型商业应用先导区”，",
    "以明确定位提升物业溢价。"], accent=BLUE)
# 产品形态条
rect(s, Inches(0.55), Inches(4.45), Inches(12.2), Inches(2.4), RGBColor(0xEA, 0xF6, 0xF8))
textbox(s, Inches(0.8), Inches(4.6), Inches(11.5), Inches(0.5),
        [[("空间产品形态（把物业变成可持续招商、可变现的“产业产品”）", 15, BLUE, True)]])
for i, (t) in enumerate(["AI 企业轻办公", "AI 应用展示厅（可招商）", "AI 路演 / 活动厅（变现）", "算力接入点（云基地结合）"]):
    x = Inches(0.8 + i * 3.0)
    rect(s, x, Inches(5.2), Inches(2.75), Inches(1.4), WHITE)
    rect(s, x, Inches(5.2), Inches(0.1), Inches(1.4), CYAN)
    textbox(s, x + Inches(0.2), Inches(5.2), Inches(2.5), Inches(1.4),
            [[(t, 14, DARKTX, True)]], anchor=MSO_ANCHOR.MIDDLE)


# ============================================================================
# 幻灯片 6 — 合作点2 招商对接 分层招商表
# ============================================================================
s = add_slide()
section_header(s, "合作点 2 · 招商对接", "从“盲目找租客”转向“生态圈定向导入”", "04")
textbox(s, Inches(0.55), Inches(1.5), Inches(12), Inches(0.5),
        [[("不是招商，而是帮对方构建一个产业生态结构 —— 分层招商模型", 15, BLUE, True)]])
rows, cols = 4, 4
tbl_shape = s.shapes.add_table(rows, cols, Inches(0.55), Inches(2.05), Inches(7.7), Inches(2.6))
table = tbl_shape.table
table.columns[0].width = Inches(1.0)
table.columns[1].width = Inches(2.3)
table.columns[2].width = Inches(2.2)
table.columns[3].width = Inches(2.2)
fill_table(table, [
    ["层级", "类型", "来源", "作用"],
    ["S 层", "AI 独角兽 / 龙头", "组委会 / 峰会资源", "提升品牌"],
    ["A 层", "成长型 AI 公司", "联合会 / 复旦", "稳定租金"],
    ["B 层", "初创团队", "学院 / 孵化器", "填充空间"],
])
style_table(table, header_fill=NAVY)
# 右侧三个增值抓手
bullet(s, Inches(8.6), Inches(2.0), Inches(4.2), Inches(4.6), [
    ("承接峰会长尾红利：", "将联合会内有扩张/选址需求企业作为首批“种子租客”定向导入。"),
    ("AI 落地赋能包作招商诱饵：", "AI 招聘/视频矩阵/数据分析等工作流打包为“智慧企业入驻包”。"),
    ("科技金融加速器节点：", "联合泰隆银行等，设专属科技金融触点，空间+技术+资本闭环。"),
], size=13, lead_color=CYAN)


# ============================================================================
# 幻灯片 7 — 杀手锏 Demo City
# ============================================================================
s = add_slide()
section_header(s, "差异化杀手锏", "AI 商业化落地试验场（Demo City）", "05")
textbox(s, Inches(0.55), Inches(1.5), Inches(12), Inches(0.6),
        [[("把这里做成“上海 AI 商业化落地样板点”，而不是一个普通产业园。", 16, BLUE, True)]])
# 三要素
for i, (t, lines, c) in enumerate([
    ("云计算基地", ["提供算力", "提供真实落地场景"], BLUE),
    ("创新创业学院", ["提供人才", "提供课程与培训"], CYAN),
    ("我们四方", ["提供项目与企业", "提供商业化方法论"], GOLD),
]):
    x = Inches(0.55 + i * 4.15)
    card(s, x, Inches(2.2), Inches(3.85), Inches(1.9), t, lines, accent=c,
         title_color=(NAVY if c == GOLD else WHITE))
# 落地场景
rect(s, Inches(0.55), Inches(4.4), Inches(12.2), Inches(2.4), RGBColor(0xEA, 0xF6, 0xF8))
textbox(s, Inches(0.8), Inches(4.55), Inches(11.5), Inches(0.5),
        [[("三大落地场景", 15, BLUE, True)]])
for i, (t, d) in enumerate([
    ("AI + 零售", "依托大学路商业场景"),
    ("AI + 商业地产", "空置物业智能盘活"),
    ("AI + 办公", "智能办公空间样板"),
]):
    x = Inches(0.8 + i * 4.0)
    rect(s, x, Inches(5.15), Inches(3.7), Inches(1.45), WHITE)
    rect(s, x, Inches(5.15), Inches(0.1), Inches(1.45), GOLD)
    textbox(s, x + Inches(0.22), Inches(5.25), Inches(3.4), Inches(1.3),
            [[(t, 16, DARKTX, True)], [(d, 12.5, GREY, False)]], space_after=Pt(4),
            anchor=MSO_ANCHOR.MIDDLE)


# ============================================================================
# 幻灯片 7.5 — 腾讯云加入：算力与技术底座
# ============================================================================
s = add_slide()
section_header(s, "新伙伴 · 腾讯云", "引入腾讯云：补齐“算力 + 技术”全栈底座", "+")
textbox(s, Inches(0.55), Inches(1.45), Inches(12), Inches(0.55),
        [[("生态从“空间 + 产业”升级为“空间 + 产业 + 算力 + 技术”全栈闭环。", 16, BLUE, True)]])
tx_cards = [
    ("① 算力补贴", ["为入驻 OPC / 企业发放", "云资源代金券与算力补贴，", "直接降本，跑通 AI 应用。"], BLUE),
    ("② 大模型与 API", ["混元大模型、API、", "Agent 开发工具，", "加速 AI 应用快速落地。"], CYAN),
    ("③ 技术中台", ["云开发、数据库、音视频", "等中台能力，", "支撑 AI 工作流规模化。"], BLUE),
    ("④ 创业扶持", ["腾讯云初创扶持计划 +", "生态资源对接，", "技术背书提升招商力。"], GOLD),
]
for i, (t, lines, c) in enumerate(tx_cards):
    x = Inches(0.55 + i * 3.08)
    card(s, x, Inches(2.1), Inches(2.85), Inches(2.55), t, lines, accent=c,
         title_color=(NAVY if c == GOLD else WHITE))
rect(s, Inches(0.55), Inches(4.95), Inches(12.2), Inches(1.85), NAVY)
rect(s, Inches(0.55), Inches(4.95), Inches(0.16), Inches(1.85), CYAN)
textbox(s, Inches(1.0), Inches(5.15), Inches(11.5), Inches(1.5),
        [[("战略意义：", 16, CYAN, True),
          ("算力是 OPC（AI 一人公司）的“水电煤”。", 16, WHITE, False)],
         [("腾讯云作底座，让“超级孵化器”能以极低成本批量孵化 OPC，把算力补贴转化为入驻黏性与招商利器，", 15, LIGHT, False)],
         [("最终形成“算力 → 孵化 → 落地 → 集聚”的产业闭环。", 15, GOLD, True)]],
        space_after=Pt(6), line_spacing=1.08)


# ============================================================================
# 幻灯片 7.6 — 超级孵化器 · OPC 孵化生态闭环（核心）
# ============================================================================
s = add_slide()
section_header(s, "CORE · 产业闭环", "超级孵化器 · OPC 孵化生态闭环", "★")
textbox(s, Inches(0.55), Inches(1.42), Inches(12.2), Inches(0.95),
        [[("OPC = AI 一人公司 / 超级个体：", 15, BLUE, True),
          ("借助 AI 工作流，以极小团队完成过去整家公司的业务。五方联合「超级孵化器」批量孵化 OPC —— 这是生态的最小细胞，也是闭环的主角。", 14, GREY, False)]],
        line_spacing=1.08)
steps = [
    ("1 · 选育", "AI 切磋者大会", ["海选 / 比武发现优质 OPC 与项目", "组委会 + 联合会输送种子"], BLUE),
    ("2 · 赋能", "超级孵化器", ["腾讯云算力补贴 + AI 工作流导入", "复旦课程 / 导师 / 政策"], CYAN),
    ("3 · 落地", "推动 AI 落地", ["Demo City 真实场景验证", "联合会企业订单 + 投融资对接"], GOLD),
    ("4 · 集聚", "入驻成链", ["OPC 入驻盘活空间，形成生态链", "反哺招商、品牌与资产增值"], BLUE),
]
for i, (no, t, lines, c) in enumerate(steps):
    x = Inches(0.55 + i * 3.02)
    rect(s, x, Inches(2.55), Inches(2.75), Inches(2.55), LIGHT)
    rect(s, x, Inches(2.55), Inches(2.75), Inches(0.95), c)
    textbox(s, x, Inches(2.55), Inches(2.75), Inches(0.5),
            [[(no, 14, (NAVY if c == GOLD else WHITE), True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, x, Inches(3.0), Inches(2.75), Inches(0.5),
            [[(t, 16, (NAVY if c == GOLD else WHITE), True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, x + Inches(0.16), Inches(3.65), Inches(2.45), Inches(1.4),
            [[(l, 11.5, DARKTX, False)] for l in lines], space_after=Pt(4), line_spacing=1.05)
    if i < 3:
        ar = rect(s, Inches(0.55 + i * 3.02 + 2.78), Inches(3.45), Inches(0.22), Inches(0.7),
                  CYAN, shape=MSO_SHAPE.CHEVRON)
# 闭环底栏
rect(s, Inches(0.55), Inches(5.4), Inches(12.2), Inches(1.4), NAVY)
rect(s, Inches(0.55), Inches(5.4), Inches(0.16), Inches(1.4), GOLD)
textbox(s, Inches(1.0), Inches(5.55), Inches(11.5), Inches(1.15),
        [[("↺  形成生态链 · 服务目前所有 OPC · 闭环自循环：", 16, GOLD, True),
          ("入驻 OPC 成长后产生新需求与新订单，又反哺孵化与招商，", 15, WHITE, False)],
         [("使空间持续满租、产业持续集聚，最终沉淀为可复制到全国的“AI 产业闭环”产品。", 15, LIGHT, False)]],
        space_after=Pt(6), line_spacing=1.08)


# ============================================================================
# 幻灯片 8 — 内容运营 / 活动 IP
# ============================================================================
s = add_slide()
section_header(s, "运营引流 · 治本之策", "内容 + 活动引流：空间是容器，没有内容必然空置", "06")
items = [
    ("AI 切磋者大会", "选育 · 招商", GOLD),
    ("AI 产业私董会", "高端 · 圈层", BLUE),
    ("AI Demo Day", "路演 · 转化", CYAN),
    ("AI 应用展", "To G / To B", BLUE),
    ("复旦+产业课程", "收费 · 人才", CYAN),
]
for i, (t, d, c) in enumerate(items):
    x = Inches(0.45 + i * 2.52)
    rect(s, x, Inches(1.9), Inches(2.3), Inches(2.3), LIGHT)
    rect(s, x, Inches(1.9), Inches(2.3), Inches(0.7), c)
    textbox(s, x, Inches(1.9), Inches(2.3), Inches(0.7),
            [[("IP %d" % (i + 1), 13, (NAVY if c == GOLD else WHITE), True)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, x + Inches(0.1), Inches(2.7), Inches(2.1), Inches(1.4),
            [[(t, 15, DARKTX, True)], [(d, 12.5, GREY, False)]], align=PP_ALIGN.CENTER,
            space_after=Pt(6))
rect(s, Inches(0.55), Inches(4.6), Inches(12.2), Inches(2.1), NAVY)
rect(s, Inches(0.55), Inches(4.6), Inches(0.16), Inches(2.1), CYAN)
textbox(s, Inches(1.0), Inches(4.85), Inches(11.5), Inches(1.6),
        [[("核心逻辑：空间 = 内容容器。其中「AI 切磋者大会」是 OPC 的选育入口，", 16, WHITE, False)],
         [("固定 IP 持续制造流量、选育 OPC、产出签约线索，从根本上解决空置并反哺招商与品牌溢价。", 16, GOLD, True)]],
        space_after=Pt(10))


# ============================================================================
# 幻灯片 9 — 四方分工与协同机制 表
# ============================================================================
s = add_slide()
section_header(s, "GOVERNANCE · 合作阵型", "五方分工与协同机制", "07")
rows, cols = 6, 3
tshape = s.shapes.add_table(rows, cols, Inches(0.55), Inches(1.5), Inches(12.2), Inches(4.45))
table = tshape.table
table.columns[0].width = Inches(3.4)
table.columns[1].width = Inches(3.6)
table.columns[2].width = Inches(5.2)
fill_table(table, [
    ["组织方", "定位 / 核心职责", "落地抓手"],
    ["云计算创新基地 / 五角场创新创业学院", "产业原点 · 孵化体系与空间品牌", "产业准入标准、开放活动平台、输出运营模式"],
    ["腾讯云", "算力与技术底座 · 全栈支撑", "算力补贴、大模型/API、技术中台、初创扶持"],
    ["复旦大学住房政策研究中心", "智库大脑 · 资产评估与政策支持", "物业改造可行性报告、城市更新补贴、人才公寓配套"],
    ["人工智能商业化落地组委会", "流量漏斗(技术) · AI 内容与投融资", "AI 切磋者大会、OPC 选育、商业化落地辅导"],
    ["杨浦区科技企业联合会", "流量漏斗(企业) · 招商渠道与跨区联动", "企业需求数据库、组织会员考察、对接各地商会"],
])
style_table(table, header_fill=NAVY, body_size=11.5)
textbox(s, Inches(0.55), Inches(6.15), Inches(12.2), Inches(1.0),
        [[("协同机制：", 13, BLUE, True), ("五方联席工作组（每月对齐）· 资源共享清单 · 战略合作备忘录（明确佣金、运营分成与算力补贴口径）。", 13, GREY, False)]])


# ============================================================================
# 幻灯片 10 — 收益模式 表
# ============================================================================
s = add_slide()
section_header(s, "BUSINESS MODEL · 收益结构", "合作分工 × 收益模式：拿到“运营权 + 导入权”", "08")
# 左：分工表
textbox(s, Inches(0.55), Inches(1.5), Inches(6), Inches(0.4), [[("合作分工", 14, BLUE, True)]])
t1 = s.shapes.add_table(6, 3, Inches(0.55), Inches(1.9), Inches(6.0), Inches(3.5)).table
t1.columns[0].width = Inches(1.3); t1.columns[1].width = Inches(2.35); t1.columns[2].width = Inches(2.35)
fill_table(t1, [
    ["模块", "平台方（云基地/腾讯云）", "我们"],
    ["空间", "提供物业", "选址 + 定位"],
    ["算力技术", "腾讯云补贴 + 中台", "OPC 孵化 + 工作流"],
    ["招商", "政策资源", "企业导入"],
    ["运营", "基础管理", "内容 + 活动 + 社群"],
    ["品牌", "政府/腾讯云背书", "AI 产业 IP"],
])
style_table(t1, header_fill=BLUE, body_size=11.5, header_size=11.5)
# 右：收益
textbox(s, Inches(7.0), Inches(1.5), Inches(6), Inches(0.4), [[("收入模式（我方主动提）", 14, CYAN, True)]])
bullet(s, Inches(7.0), Inches(1.9), Inches(5.8), Inches(3.5), [
    ("① 租金分成", "—— 核心收益来源。"),
    ("② 企业服务费", "—— 培训 / 咨询 / AI 赋能包。"),
    ("③ 活动收入", "—— 切磋者大会 / Demo Day / 展会。"),
    ("④ 算力 / 技术服务分成", "—— 腾讯云算力分销 + OPC 孵化服务。"),
    ("⑤ 政策补贴联合申报", "—— 城市更新、保租房、算力补贴。"),
], size=13, lead_color=CYAN)
rect(s, Inches(0.55), Inches(5.7), Inches(12.2), Inches(1.3), RGBColor(0xEA, 0xF6, 0xF8))
rect(s, Inches(0.55), Inches(5.7), Inches(0.16), Inches(1.3), GOLD)
textbox(s, Inches(0.95), Inches(5.82), Inches(11.6), Inches(1.05),
        [[("关键原则：", 15, BLUE, True),
          ("我方必须拿到“运营权 + 产业导入权”，否则只是给资源平台打工，无法形成可复制的产品与护城河。", 15, DARKTX, False)]],
        anchor=MSO_ANCHOR.MIDDLE)


# ============================================================================
# 幻灯片 11 — 三个立即启动项目
# ============================================================================
s = add_slide()
section_header(s, "ACTION · 即刻启动", "四个“立即启动”项目（承接 5·22 峰会余热）", "09")
projects = [
    ("01", "AI 产业社区试点", "联合考察对方推荐的 1–2 个空置物业（大学路 / 杨浦滨江）。",
     "交付：2 周内出具《产业社区改造初步方案》", BLUE),
    ("02", "首场 AI 切磋者大会 + OPC 超级孵化营", "承接峰会余热，7 月办首站，选育首批种子 OPC。",
     "交付：活动方案 + 拟邀 OPC/企业名单（本周完成）", GOLD),
    ("03", "腾讯云算力补贴包落地", "联合腾讯云敲定面向入驻 OPC 的算力补贴与扶持政策。",
     "交付：2 周内确定补贴口径、额度与申领流程", CYAN),
    ("04", "企业需求数据库共建", "本周启动联合会会员企业扩租 / 迁址需求调研。",
     "交付：1 个月内形成首批 50 家目标企业清单", BLUE),
]
for i, (no, t, d, deliver, c) in enumerate(projects):
    y = Inches(1.55 + i * 1.38)
    rect(s, Inches(0.55), y, Inches(12.2), Inches(1.22), LIGHT)
    rect(s, Inches(0.55), y, Inches(1.3), Inches(1.22), c)
    textbox(s, Inches(0.55), y, Inches(1.3), Inches(1.22),
            [[(no, 30, (NAVY if c == GOLD else WHITE), True)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(2.05), y + Inches(0.08), Inches(10.4), Inches(1.1),
            [[(t, 16, DARKTX, True)],
             [(d, 12.5, GREY, False)],
             [(deliver, 12.5, BLUE, True)]], space_after=Pt(2), line_spacing=1.03)


# ============================================================================
# 幻灯片 12 — 价值主张 + 收口句
# ============================================================================
s = add_slide()
set_bg(s, NAVY)
rect(s, 0, 0, Inches(0.22), EMU_H, CYAN)
rect(s, Inches(0.22), 0, Inches(0.08), EMU_H, GOLD)
textbox(s, Inches(0.9), Inches(0.85), Inches(11.5), Inches(0.5),
        [[("核心价值主张", 16, CYAN, True)]])
textbox(s, Inches(0.9), Inches(1.4), Inches(11.6), Inches(1.6),
        [[("“我们不是来租您房子的 ——", 26, WHITE, True)],
         [("  是来帮您把空置物业变成 AI 产业资产的。”", 26, GOLD, True)]],
        space_after=Pt(8), line_spacing=1.1)
# 三个不一样
for i, (t, d) in enumerate([
    ("模式不一样", "不做二房东，做产业运营服务商"),
    ("招商不一样", "不靠中介，靠产业生态精准匹配"),
    ("价值不一样", "不只填空置，更提升资产估值"),
]):
    x = Inches(0.9 + i * 4.0)
    rect(s, x, Inches(3.2), Inches(3.7), Inches(1.5), RGBColor(0x12, 0x33, 0x5E))
    rect(s, x, Inches(3.2), Inches(3.7), Inches(0.08), CYAN)
    textbox(s, x + Inches(0.2), Inches(3.35), Inches(3.4), Inches(1.3),
            [[(t, 16, GOLD, True)], [(d, 13, LIGHT, False)]], space_after=Pt(6))
# 收口句
rect(s, Inches(0.9), Inches(5.1), Inches(11.5), Inches(1.7), RGBColor(0x12, 0x33, 0x5E))
rect(s, Inches(0.9), Inches(5.1), Inches(0.16), Inches(1.7), GOLD)
textbox(s, Inches(1.3), Inches(5.25), Inches(11), Inches(1.4),
        [[("收口句：", 15, CYAN, True)],
         [("“如果今天只是谈招商，我们做的是一单生意；", 19, WHITE, True)],
         [("  如果我们一起做模型，这是一个可以复制到全国的产品。”", 19, WHITE, True)]],
        space_after=Pt(4), line_spacing=1.1)


prs.save("合作方案_AI云计算产业生态共同体.pptx")
print("OK -> 合作方案_AI云计算产业生态共同体.pptx  slides=%d" % len(prs.slides._sldIdLst))
