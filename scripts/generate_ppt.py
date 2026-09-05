# -*- coding: utf-8 -*-
"""生成《外滩·产业课堂》联合策划方案 PPT。"""

from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

import content as C

OUT = Path(__file__).resolve().parent.parent / "output" / "港大经管上海中心_外滩产业课堂_联合策划方案.pptx"

# 港大经管常见绿金，便于对方打开时有熟悉感
GREEN = RGBColor(0x00, 0x3D, 0x2E)
GREEN2 = RGBColor(0x0A, 0x5C, 0x46)
GOLD = RGBColor(0xC4, 0xA3, 0x5A)
CREAM = RGBColor(0xF7, 0xF4, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x24, 0x20)
GREY = RGBColor(0x5B, 0x6B, 0x64)
LIGHT = RGBColor(0xE8, 0xEF, 0xEA)

FONT = "微软雅黑"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL = 13


def set_font(run, size, color=DARK, bold=False, font=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = rpr.makeelement(qn("a:ea"), {})
        rpr.append(ea)
    ea.set("typeface", font)


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_round(slide, x, y, w, h, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    # 略收圆角
    try:
        shp.adjustments[0] = 0.08
    except Exception:
        pass
    return shp


def add_text(slide, x, y, w, h, text, size, color=DARK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_font(r, size, color, bold)
    return tb


def add_runs(slide, x, y, w, h, paragraphs, anchor=MSO_ANCHOR.TOP):
    """paragraphs: list of list of (text, size, color, bold)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    first = True
    for segs in paragraphs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(6)
        for text, size, color, bold in segs:
            r = p.add_run()
            r.text = text
            set_font(r, size, color, bold)
    return tb


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def header(slide, title, idx):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.05), GREEN)
    add_rect(slide, 0, Inches(1.05), SLIDE_W, Pt(3.5), GOLD)
    add_text(slide, Inches(0.55), Inches(0.18), Inches(11.2), Inches(0.7), title, 24, WHITE, True, anchor=MSO_ANCHOR.MIDDLE)
    add_text(
        slide,
        Inches(11.5),
        Inches(0.18),
        Inches(1.4),
        Inches(0.7),
        f"{idx:02d}/{TOTAL:02d}",
        12,
        GOLD,
        True,
        align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def footer(slide):
    add_text(
        slide,
        Inches(0.55),
        Inches(7.12),
        Inches(12.2),
        Inches(0.28),
        f"{C.PROJECT_NAME}  ·  {C.THEIR_UNIT} × {C.OUR_PARTIES}  ·  {C.VERSION}  ·  {C.DATE_CN}",
        9,
        GREY,
    )


def bullets(slide, x, y, w, h, items, size=15, color=DARK, gap=8):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.line_spacing = 1.15
        rm = p.add_run()
        rm.text = "●  "
        set_font(rm, size - 2, GOLD, True)
        r = p.add_run()
        r.text = it
        set_font(r, size, color)
    return tb


# ---------------------------------------------------------------------------
# 各页
# ---------------------------------------------------------------------------


def slide_cover(prs):
    s = blank(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, GREEN)
    add_rect(s, 0, 0, Inches(0.18), SLIDE_H, GOLD)
    add_text(s, Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.35), C.THEIR_MOTTO, 14, GOLD, True)
    add_text(s, Inches(0.8), Inches(1.15), Inches(11.5), Inches(0.4), C.THEIR_UNIT, 16, WHITE, False)
    add_text(s, Inches(0.8), Inches(2.15), Inches(11.8), Inches(0.9), C.PROJECT_NAME, 44, WHITE, True)
    add_text(s, Inches(0.8), Inches(3.05), Inches(11.8), Inches(0.5), C.PROJECT_SUBTITLE, 22, GOLD, False)
    add_rect(s, Inches(0.8), Inches(3.75), Inches(2.2), Pt(3), GOLD)
    add_text(
        s,
        Inches(0.8),
        Inches(4.05),
        Inches(11.5),
        Inches(1.2),
        C.ONE_LINER,
        16,
        WHITE,
        False,
    )
    add_text(
        s,
        Inches(0.8),
        Inches(5.7),
        Inches(11.5),
        Inches(0.7),
        f"致：{C.THEIR_CONTACT}  {C.THEIR_TITLE}\n{C.OUR_PARTIES}  |  {C.VERSION}  |  {C.DATE_CN}",
        14,
        RGBColor(0xD5, 0xE4, 0xDC),
        False,
    )
    add_text(s, Inches(0.8), Inches(6.85), Inches(11.5), Inches(0.35), C.THEIR_ADDR + "  ·  " + C.THEIR_WEB, 12, GOLD, False)


def slide_heard(prs):
    s = blank(prs)
    header(s, f"{C.MEETING_DATE} 交流：先听懂对方怎么说", 2)
    add_round(s, Inches(0.5), Inches(1.28), Inches(12.3), Inches(1.05), CREAM)
    add_text(
        s,
        Inches(0.75),
        Inches(1.4),
        Inches(11.8),
        Inches(0.8),
        C.ONE_LINER,
        16,
        DARK,
        False,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    for i, (t, d) in enumerate(C.HEARD):
        x = Inches(0.5 + (i % 2) * 6.25)
        y = Inches(2.55 + (i // 2) * 2.05)
        add_round(s, x, y, Inches(6.05), Inches(1.9), WHITE)
        add_rect(s, x, y, Inches(0.12), Inches(1.9), GOLD)
        add_text(s, x + Inches(0.35), y + Inches(0.18), Inches(5.5), Inches(0.4), t, 16, GREEN, True)
        add_text(s, x + Inches(0.35), y + Inches(0.62), Inches(5.5), Inches(1.1), d, 14, DARK, False)
    footer(s)


def slide_needs(prs):
    s = blank(prs)
    header(s, "对方需求四件事（按洽谈整理，不是我们替对方编的）", 3)
    for i, (t, d) in enumerate(C.THEIR_NEEDS):
        y = Inches(1.32 + i * 1.32)
        add_round(s, Inches(0.5), y, Inches(12.3), Inches(1.18), WHITE)
        add_rect(s, Inches(0.5), y, Inches(0.14), Inches(1.18), GOLD)
        add_text(s, Inches(0.85), y + Inches(0.12), Inches(11.7), Inches(0.38), f"{i + 1}.  {t}", 16, GREEN, True)
        add_text(s, Inches(0.85), y + Inches(0.52), Inches(11.7), Inches(0.55), d, 13, GREY, False)
    footer(s)


def slide_window(prs):
    s = blank(prs)
    header(s, "上海中心怎么运转：四个点位，四项职能", 4)
    for i, (name, desc) in enumerate(C.CENTRE_SITES):
        x = Inches(0.45 + i * 3.2)
        add_round(s, x, Inches(1.32), Inches(3.05), Inches(2.35), WHITE)
        add_rect(s, x, Inches(1.32), Inches(3.05), Inches(0.55), GREEN if i == 1 else GREEN2)
        add_text(s, x, Inches(1.32), Inches(3.05), Inches(0.55), name, 14, WHITE, True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.12), Inches(2.0), Inches(2.8), Inches(1.5), desc, 12, DARK, False, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.55), Inches(3.85), Inches(12.2), Inches(0.35), "经管中心当前职能（对方口述）", 16, GREEN, True)
    bullets(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.5), C.CENTRE_FUNCTIONS, size=14, gap=6)
    footer(s)


def slide_gap(prs):
    s = blank(prs)
    header(s, "双方互补：出海是对方点名的赛道", 5)
    labels = C.GAP
    colors = [GREEN, GREEN2, GREEN, GREEN2]
    for i, ((k, v), col) in enumerate(zip(labels, colors)):
        x = Inches(0.5 + i * 3.15)
        add_round(s, x, Inches(1.4), Inches(3.0), Inches(3.35), WHITE)
        add_rect(s, x, Inches(1.4), Inches(3.0), Inches(0.7), col)
        add_text(s, x, Inches(1.4), Inches(3.0), Inches(0.7), k, 16, WHITE, True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.18), Inches(2.3), Inches(2.65), Inches(2.2), v, 14, DARK, False, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.55), Inches(4.95), Inches(12.2), Inches(0.4), "我方可以立刻用上的资产", 16, GREEN, True)
    bullets(s, Inches(0.55), Inches(5.32), Inches(12.2), Inches(1.7), C.OUR_ASSETS, size=12, gap=3)
    footer(s)


def slide_product(prs):
    s = blank(prs)
    header(s, "一条思路：三层合作，出海是主赛道", 6)
    add_text(s, Inches(0.55), Inches(1.28), Inches(12.2), Inches(0.32), "按 9 月 3 日对方口径收口", 14, GREEN, True)
    bullets(s, Inches(0.55), Inches(1.6), Inches(12.2), Inches(1.45), C.PRODUCT_LOGIC, size=14, gap=4)
    for i, p in enumerate(C.LAYERS):
        x = Inches(0.45 + i * 4.2)
        add_round(s, x, Inches(3.2), Inches(4.05), Inches(3.2), WHITE)
        add_rect(s, x, Inches(3.2), Inches(4.05), Inches(0.12), GOLD)
        add_text(s, x + Inches(0.18), Inches(3.42), Inches(3.7), Inches(0.7), p["name"], 15, GREEN, True)
        add_text(s, x + Inches(0.18), Inches(4.15), Inches(3.7), Inches(0.32), p["tag"], 11, GOLD, True)
        add_text(s, x + Inches(0.18), Inches(4.52), Inches(3.7), Inches(1.65), p["desc"], 13, DARK, False)
    footer(s)


def slide_convert(prs):
    s = blank(prs)
    header(s, "如何变成生源：现场谈出海，会后由招生官收口", 7)
    steps = [
        ("1", "邀约", "按 EMBA 门槛筛名单\n企业主 / 高管优先"),
        ("2", "到场", "出海议题把人留下\n港大路径把场撑住"),
        ("3", "分级", "会后 48h 分 A/B/C\n名单共管不外流"),
        ("4", "跟进", "A 类潘老师约面\n录取权始终在学院"),
    ]
    for i, (n, t, d) in enumerate(steps):
        x = Inches(0.55 + i * 3.2)
        add_round(s, x, Inches(1.45), Inches(2.95), Inches(2.7), WHITE)
        add_round(s, x + Inches(1.1), Inches(1.65), Inches(0.7), Inches(0.7), GOLD)
        add_text(s, x + Inches(1.1), Inches(1.65), Inches(0.7), Inches(0.7), n, 20, GREEN, True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, Inches(2.45), Inches(2.95), Inches(0.4), t, 18, GREEN, True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.15), Inches(2.95), Inches(2.65), Inches(1.0), d, 13, GREY, False, align=PP_ALIGN.CENTER)
        if i < 3:
            add_text(s, x + Inches(2.7), Inches(2.4), Inches(0.5), Inches(0.4), "→", 20, GOLD, True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.55), Inches(4.4), Inches(12.2), Inches(0.4), "首场样板 KPI（用于验收前期费用，不是录取承诺）", 16, GREEN, True)

    # 简易表
    headers = ["指标", "约定"]
    rows = C.KPI
    table = s.shapes.add_table(len(rows) + 1, 2, Inches(0.55), Inches(4.85), Inches(12.2), Inches(1.9)).table
    table.columns[0].width = Inches(2.4)
    table.columns[1].width = Inches(9.8)
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = h
        set_font(r, 13, WHITE, True)
    for ri, (a, b) in enumerate(rows, 1):
        for ci, val in enumerate((a, b)):
            cell = table.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci == 0 else PP_ALIGN.LEFT
            r = p.add_run()
            r.text = val
            set_font(r, 12, DARK, ci == 0)
    footer(s)


def slide_90(prs):
    s = blank(prs)
    header(s, "90 天：外滩首场 + 一次主任体验", 8)
    for i, (when, title, desc) in enumerate(C.NINETY_DAY):
        y = Inches(1.35 + i * 1.3)
        add_round(s, Inches(0.5), y, Inches(12.3), Inches(1.18), WHITE)
        add_rect(s, Inches(0.5), y, Inches(0.14), Inches(1.18), GOLD)
        add_text(s, Inches(0.85), y + Inches(0.18), Inches(2.4), Inches(0.8), when, 16, GREEN, True, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(3.3), y + Inches(0.15), Inches(9.2), Inches(0.4), title, 16, DARK, True)
        add_text(s, Inches(3.3), y + Inches(0.55), Inches(9.2), Inches(0.5), desc, 13, GREY, False)
    footer(s)


def slide_first(prs):
    s = blank(prs)
    header(s, "首场怎么开：出海主题可改，赛道不改", 9)
    ev = C.FIRST_EVENT
    add_round(s, Inches(0.5), Inches(1.3), Inches(5.5), Inches(5.5), WHITE)
    add_text(s, Inches(0.75), Inches(1.5), Inches(5.1), Inches(0.4), ev["title"], 16, GOLD, True)
    meta = [
        ("主题", ev["theme"]),
        ("时间", ev["when"]),
        ("地点", ev["where"]),
        ("人群", ev["who"]),
    ]
    yy = 2.05
    for k, v in meta:
        add_text(s, Inches(0.75), Inches(yy), Inches(5.1), Inches(0.3), k, 12, GREEN, True)
        add_text(s, Inches(0.75), Inches(yy + 0.28), Inches(5.1), Inches(0.7), v, 13, DARK, False)
        yy += 1.05

    add_round(s, Inches(6.2), Inches(1.3), Inches(6.6), Inches(5.5), CREAM)
    add_text(s, Inches(6.45), Inches(1.5), Inches(6.2), Inches(0.4), "半日议程", 16, GREEN, True)
    bullets(s, Inches(6.45), Inches(2.0), Inches(6.15), Inches(4.5), ev["agenda"], size=13, gap=6)
    footer(s)


def slide_roles(prs):
    s = blank(prs)
    header(s, "分工：招生官做粘性，主任入口另走一层", 10)
    keys = list(C.ROLES.keys())
    for i, k in enumerate(keys):
        x = Inches(0.5 + i * 6.4)
        add_round(s, x, Inches(1.4), Inches(6.15), Inches(5.35), WHITE)
        add_rect(s, x, Inches(1.4), Inches(6.15), Inches(0.7), GREEN if i == 0 else GREEN2)
        add_text(s, x, Inches(1.4), Inches(6.15), Inches(0.7), k, 18, WHITE, True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        bullets(s, x + Inches(0.3), Inches(2.3), Inches(5.55), Inches(4.1), C.ROLES[k], size=15, gap=10)
    footer(s)


def slide_fee(prs):
    s = blank(prs)
    header(s, "商业结构：只收一笔前期费用", 11)
    add_round(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.55), GREEN)
    add_text(s, Inches(0.75), Inches(1.42), Inches(12.0), Inches(0.35), C.FEE_NAME, 14, GOLD, True)
    add_text(
        s,
        Inches(0.75),
        Inches(1.78),
        Inches(12.0),
        Inches(0.85),
        f"{C.FEE_AMOUNT_CN}    ¥{C.FEE_AMOUNT:,}\n合同生效后 {C.FEE_DAYS} 个工作日内一次性支付，开具增值税发票",
        22,
        WHITE,
        True,
    )

    add_text(s, Inches(0.55), Inches(3.05), Inches(6.0), Inches(0.4), "这 8.8 万覆盖什么", 16, GREEN, True)
    bullets(s, Inches(0.5), Inches(3.4), Inches(6.1), Inches(3.3), C.FEE_COVERS, size=13, gap=5)

    add_text(s, Inches(6.85), Inches(3.05), Inches(6.0), Inches(0.4), "故意不写进本协议的", 16, GREEN, True)
    bullets(s, Inches(6.8), Inches(3.4), Inches(6.1), Inches(3.3), C.FEE_NOT_COVERED, size=13, gap=5)
    footer(s)


def slide_agreement(prs):
    s = blank(prs)
    header(s, "协议要点：短、可签，按对方口径分层验收", 12)
    bullets(s, Inches(0.55), Inches(1.32), Inches(12.2), Inches(4.55), C.AGREEMENT_POINTS, size=14, gap=7)
    add_round(s, Inches(0.5), Inches(5.95), Inches(12.3), Inches(0.9), CREAM)
    add_text(
        s,
        Inches(0.75),
        Inches(6.05),
        Inches(11.8),
        Inches(0.7),
        "完整文本见配套《联合策划服务合作协议（建议稿）》。本稿供潘老师内部预审，签署主体、账号、开票信息在签署页一次性补齐。",
        14,
        DARK,
        False,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    footer(s)


def slide_next(prs):
    s = blank(prs)
    header(s, "下一步：三件事，合作就算成立", 13)
    for i, t in enumerate(C.NEXT_STEPS):
        y = Inches(1.4 + i * 1.35)
        add_round(s, Inches(0.55), y, Inches(12.2), Inches(1.2), WHITE)
        add_round(s, Inches(0.8), y + Inches(0.28), Inches(0.65), Inches(0.65), GOLD)
        add_text(s, Inches(0.8), y + Inches(0.28), Inches(0.65), Inches(0.65), str(i + 1), 20, GREEN, True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.7), y + Inches(0.25), Inches(10.7), Inches(0.7), t, 16, DARK, False, anchor=MSO_ANCHOR.MIDDLE)
    add_text(
        s,
        Inches(0.55),
        Inches(5.6),
        Inches(12.2),
        Inches(1.2),
        f"接口人建议：{C.THEIR_CONTACT}  {C.THEIR_TITLE}\n{C.THEIR_TEL}  ·  {C.THEIR_EMAIL}  ·  {C.THEIR_ADDR}",
        14,
        GREY,
        False,
        align=PP_ALIGN.CENTER,
    )
    footer(s)


def build(path: Path | None = None) -> Path:
    path = path or OUT
    path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide_cover(prs)
    slide_heard(prs)
    slide_needs(prs)
    slide_window(prs)
    slide_gap(prs)
    slide_product(prs)
    slide_convert(prs)
    slide_90(prs)
    slide_first(prs)
    slide_roles(prs)
    slide_fee(prs)
    slide_agreement(prs)
    slide_next(prs)
    prs.save(path)
    return path


if __name__ == "__main__":
    out = Path(sys.argv[1]) if len(sys.argv) > 1 else OUT
    p = build(out)
    print(f"已生成 {p}")
