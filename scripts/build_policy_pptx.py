"""
冠松 01# 楼 · 政策路径董事会版 PPT（8 页）

9 月 2 日共识：政府公关为主、招商为辅。
不讲 iDrive Hub，不承诺 C6 能直接做酒店。

重新生成：python3 scripts/build_policy_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pathlib import Path

NAVY = RGBColor(0x1B, 0x35, 0x55)
GOLD = RGBColor(0xB7, 0x86, 0x2E)
INK = RGBColor(0x2A, 0x2D, 0x34)
CLOUD = RGBColor(0xF0, 0xF1, 0xF3)
GREY = RGBColor(0x6B, 0x73, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0x8B, 0x2E, 0x2E)
MOSS = RGBColor(0x2E, 0x5A, 0x3C)

CN_FONT = "WenQuanYi Micro Hei"
EN_FONT = "Georgia"
SLIDES = []


def set_run(run, text, *, size=14, bold=False, color=INK, italic=False):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = EN_FONT
    rPr = run._r.get_or_add_rPr()
    for tag in ("ea", "cs"):
        existing = rPr.find(qn(f"a:{tag}"))
        if existing is not None:
            rPr.remove(existing)
        rPr.append(rPr.makeelement(qn(f"a:{tag}"), {"typeface": CN_FONT}))


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    box.fill.background()
    box.line.fill.background()
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        set_run(run, ln, size=size, bold=bold, color=color, italic=italic)
    return box


def add_rect(slide, x, y, w, h, *, fill=NAVY):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    shp.text_frame.text = ""
    return shp


def add_round(slide, x, y, w, h, text="", *, fill=NAVY, color=WHITE,
              size=12, bold=True, align=PP_ALIGN.CENTER):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p.add_run(), text, size=size, bold=bold, color=color)
    return shp


def add_table(slide, x, y, w, h, header, rows, *, header_size=12, body_size=12,
              col_widths=None):
    ts = slide.shapes.add_table(len(rows) + 1, len(header), x, y, w, h)
    table = ts.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw
    for j, htxt in enumerate(header):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.text = ""
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p.add_run(), htxt, size=header_size, bold=True, color=WHITE)
    for i, row in enumerate(rows, start=1):
        bg = WHITE if i % 2 else CLOUD
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            cell.text = ""
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            set_run(p.add_run(), str(val), size=body_size, color=INK)
    return ts


def add_chrome(slide, prs, *, page_no, label, title, subtitle=""):
    sw, sh = prs.slide_width, prs.slide_height
    add_rect(slide, 0, 0, sw, Emu(380000), fill=NAVY)
    add_round(slide, Inches(0.45), Inches(0.18), Inches(2.4), Inches(0.36),
              label, fill=GOLD, color=NAVY, size=11)
    add_text(slide, Inches(3.05), Inches(0.10), Inches(9.6), Inches(0.52),
             title, size=22, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(3.05), Inches(0.52), Inches(9.6), Inches(0.28),
                 subtitle, size=12, color=CLOUD, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 0, sh - Emu(80000), sw, Emu(80000), fill=GOLD)
    add_text(slide, Inches(0.45), Inches(7.05), Inches(9.2), Inches(0.28),
             "冠松 01# 楼 · 政策路径优化方案 · 保密", size=10, color=GREY)
    add_text(slide, Inches(11.4), Inches(7.05), Inches(1.5), Inches(0.28),
             f"{page_no} / 0", size=10, color=GREY, align=PP_ALIGN.RIGHT)


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    SLIDES.append(s)
    return s


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    sw, sh = prs.slide_width, prs.slide_height

    # 1 封面
    s = new_slide(prs)
    add_rect(s, 0, 0, sw, sh, fill=NAVY)
    add_rect(s, 0, Inches(5.05), sw, Emu(28000), fill=GOLD)
    add_text(s, Inches(0.8), Inches(1.25), Inches(11.5), Inches(0.4),
             "01# 研发楼 · 2026-09-03 优化方案",
             size=16, bold=True, color=GOLD, italic=True)
    add_text(s, Inches(0.8), Inches(1.75), Inches(11.7), Inches(1.5),
             "政府公关为主，招商为辅",
             size=36, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.4), Inches(11.7), Inches(0.9),
             "顺着董事长的公寓 / 酒店方向走\n把土地、规划、消防问题拆开，带着去问政府",
             size=18, color=CLOUD)
    add_round(s, Inches(0.8), Inches(5.5), Inches(4.2), Inches(0.48),
              "董事会版 · 8 页", fill=GOLD, color=NAVY, size=14)
    add_text(s, Inches(5.2), Inches(5.5), Inches(7.2), Inches(0.48),
             "C6 不能直接改酒店 · 9 月 30 日按已批图纸进场",
             size=14, color=CLOUD, anchor=MSO_ANCHOR.MIDDLE)

    # 2 听懂了
    s = new_slide(prs)
    add_chrome(s, prs, page_no=2, label="对齐",
               title="我们听懂了",
               subtitle="9 月 2 日沟通口径。不再劝回到纯写字楼。")
    items = [
        ("1", "不信空置写字楼",
         "金桥、陆家嘴、杨浦是董事长自己走的。诺富特满房，才是他信的对标。"),
        ("2", "要公寓 / 酒店路径",
         "一人公司、商住、特色酒店。要的是政策怎么过，不是再听一版研发招商。"),
        ("3", "9 月 30 日必须进场",
         "逾期按总造价滞纳金。第一优先是节点，不是招租名单。"),
        ("4", "要实证，过董事会",
         "照片、视频、带看、静安本地案例。口头 PPT 过不了董事长。"),
    ]
    for i, (n, t, b) in enumerate(items):
        y = Inches(1.25) + Inches(1.3) * i
        add_round(s, Inches(0.5), y, Inches(0.7), Inches(1.05), n,
                  fill=NAVY, color=WHITE, size=22)
        add_text(s, Inches(1.45), y, Inches(3.5), Inches(1.05), t,
                 size=20, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(5.1), y, Inches(7.6), Inches(1.05), b,
                 size=16, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # 3 新方向
    s = new_slide(prs)
    add_chrome(s, prs, page_no=3, label="转向",
               title="这次谈成的新方向",
               subtitle="与其把老板的头从西拧到东，不如顺着往西走，路上的问题我们拆掉。")
    boxes = [
        ("以前", NAVY,
         "独家招商填研发楼\n90 天建库带看签意向\n劝董事长相信办公能租出去"),
        ("现在", GOLD,
         "政府公关为主，招商为辅\n90 天把能不能做写成可表决材料\n公寓 / 酒店是叙事，C6 是约束"),
    ]
    for i, (t, c, b) in enumerate(boxes):
        x = Inches(0.5) + Inches(6.4) * i
        add_rect(s, x, Inches(1.3), Inches(6.05), Inches(3.35), fill=c)
        add_text(s, x, Inches(1.45), Inches(6.05), Inches(0.5), t,
                 size=18, bold=True, color=WHITE if i == 0 else NAVY,
                 align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.25), Inches(2.1), Inches(5.55), Inches(2.3),
                 b, size=18, color=WHITE if i == 0 else NAVY, align=PP_ALIGN.CENTER)

    add_rect(s, Inches(0.5), Inches(4.85), Inches(12.35), Inches(1.75), fill=CLOUD)
    add_text(s, Inches(0.75), Inches(5.0), Inches(11.9), Inches(1.45),
             "角色：政策突破操盘。复旦住房政策研究中心可出第三方身份。\n"
             "国家级孵化器品牌可叠加，授权另收费。公寓品牌可接触，没有土地路径不签约。\n"
             "无委托协议，不开展代看——没协议的对接是权责模糊期。",
             size=15, color=INK)

    # 4 红线
    s = new_slide(prs)
    add_chrome(s, prs, page_no=4, label="红线",
               title="四件事不能写进承诺",
               subtitle="可以说帮你问。不能说已经能做。")
    lines = [
        ("C6 直接做酒店 / 公寓", "教育科研设计用地，工业属性，亩产税收要求最高。硬上不合规。"),
        ("改性质几乎必然补地价", "会上量级：商业大约是工业的 10 倍。精确数以估价和规土为准。"),
        ("规划过了 ≠ 消防能开", "杨浦亚朵四号楼（会上）：保证金搁置。同项目不同楼结果可以相反。"),
        ("为等批文耽误 9/30", "施工按已批研发楼进场。业态路径并行，不拿滞纳金赌博。"),
    ]
    for i, (t, b) in enumerate(lines):
        y = Inches(1.25) + Inches(1.3) * i
        add_rect(s, Inches(0.45), y, Inches(12.4), Inches(1.18), fill=CLOUD)
        add_round(s, Inches(0.6), y + Inches(0.28), Inches(0.55), Inches(0.62),
                  str(i + 1), fill=RED, color=WHITE, size=16)
        add_text(s, Inches(1.4), y + Inches(0.08), Inches(11.2), Inches(0.5),
                 t, size=18, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.4), y + Inches(0.55), Inches(11.2), Inches(0.5),
                 b, size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # 5 三条路径
    s = new_slide(prs)
    add_chrome(s, prs, page_no=5, label="路径",
               title="三条可问的路，不是三条已批的路",
               subtitle="绿保底 · 黄主攻 · 红只做账。染色以土地证和控规图则为准。")
    paths = [
        (MOSS, "A  绿 · 保底",
         "维持 C6\n研发为主\n只做合规配套",
         "合法、不补大额地价、\n不误 9/30。\n「能住」只问人才公寓\n是否允许作科研配套。"),
        (GOLD, "B  黄 · 主攻",
         "纳入园区更新\n争取功能兼容\n不先改性质",
         "静安 5 号文第 5 条写的是\n存量商务楼宇，默认套不上。\n第 7 条园区混合利用\n才是 01# 该问的门。"),
        (RED, "C  红 · 只做账",
         "改商业 / 收储\nR 类擦边",
         "补地价、周期长。\n塘南村人才公寓出现在\n收储再出让里。\nR0 只叠在商办上。"),
    ]
    for i, (c, t, m, b) in enumerate(paths):
        x = Inches(0.4) + Inches(4.25) * i
        add_rect(s, x, Inches(1.25), Inches(4.05), Inches(5.35), fill=CLOUD)
        add_rect(s, x, Inches(1.25), Inches(4.05), Inches(0.55), fill=c)
        add_text(s, x, Inches(1.25), Inches(4.05), Inches(0.55), t,
                 size=16, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.15), Inches(1.95), Inches(3.75), Inches(1.7),
                 m, size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.15), Inches(3.7), Inches(3.75), Inches(2.6),
                 b, size=14, color=INK, align=PP_ALIGN.CENTER)

    # 6 静安案例
    s = new_slide(prs)
    add_chrome(s, prs, page_no=6, label="静安",
               title="本地窗口：有复合，没有 C6 改酒店",
               subtitle="详细来源见「14-静安研发用地业态调整-案例与政策窗口」。")
    add_table(
        s, Inches(0.4), Inches(1.25), Inches(12.5), Inches(5.4),
        ["对象", "结果", "和 01# 的距离"],
        [
            ["5 号文第 5 条", "商务楼宇可兼容酒店、人才公寓", "01# 是 C6 新建研发楼，默认不是商务楼宇"],
            ["5 号文第 7 / 13 条", "园区可混合利用；严禁改住宅", "问配套可以；商住两用易撞墙"],
            ["市规土 449 号 M0 / R0", "工研仓互转；R0 只叠商办", "C6 走 M0 到不了酒店"],
            ["走马塘 · 信谊药厂", "工业升级成研发，补地价开工", "静安愿帮升级研发，不帮降成酒店"],
            ["塘南村北块", "收储带方案，科研 + 人才公寓", "公寓在收地后再出让的方案里"],
            ["杨浦亚朵（会上）", "四号楼规划消防搁置", "外区反面：先装修不能当批准"],
        ],
        header_size=13, body_size=12,
        col_widths=[Inches(3.1), Inches(4.4), Inches(5.0)],
    )

    # 7 90天
    s = new_slide(prs)
    add_chrome(s, prs, page_no=7, label="90 天",
               title="施工照旧，政策并行，招商备着",
               subtitle="完成标准：董事会能表决的路径备忘录 + 至少一轮政府预沟通纪要。")
    phases = [
        ("D1–15 吃透地",
         "土地证、出让合同、\n控规图则、产业准入\n四条红线一页纸"),
        ("D16–45 政府对表",
         "街道 + 规资 / 更新办\n甲方出面、我方备材料\n路径 A/B/C 染色"),
        ("D46–90 可表决",
         "业态路径备忘录定稿\n补地价量级、对 9/30 影响\n招商库备着不对外承诺"),
    ]
    for i, (t, b) in enumerate(phases):
        x = Inches(0.45) + Inches(4.2) * i
        add_rect(s, x, Inches(1.25), Inches(4.0), Inches(3.7), fill=CLOUD)
        add_rect(s, x, Inches(1.25), Inches(4.0), Inches(0.6), fill=NAVY)
        add_text(s, x, Inches(1.25), Inches(4.0), Inches(0.6), t,
                 size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.15), Inches(2.05), Inches(3.7), Inches(2.7),
                 b, size=16, color=INK, align=PP_ALIGN.CENTER)

    add_rect(s, Inches(0.45), Inches(5.15), Inches(12.4), Inches(1.45), fill=NAVY)
    add_text(s, Inches(0.7), Inches(5.3), Inches(12.0), Inches(1.2),
             "复旦住房政策研究中心可出非正式意见（冲刺）。孵化器品牌授权另收费。\n"
             "无协议不代看。公寓品牌只接触。指定企业入驻、政府一定批准，这两句不写。",
             size=15, color=WHITE)

    # 8 请定
    s = new_slide(prs)
    add_chrome(s, prs, page_no=8, label="拍板",
               title="请董事会定这四件",
               subtitle="收费已在 9 月 2 日对齐，不是本页争议点。")
    asks = [
        ("01", "施工", "9 月 30 日按已批研发楼进场。业态路径不绑死开竣工。"),
        ("02", "主攻", "90 天主攻路径 B（更新/兼容）。A 保底，C 只做账。"),
        ("03", "材料", "3 个工作日内给土地证和出让合同用途全文。"),
        ("04", "身份", "是否授权复旦住房中心出面；孵化器品牌是否需要（另收费）。"),
    ]
    for i, (n, k, v) in enumerate(asks):
        y = Inches(1.22) + Inches(0.95) * i
        add_round(s, Inches(0.5), y, Inches(1.05), Inches(0.82), n,
                  fill=GOLD, color=NAVY, size=18)
        add_text(s, Inches(1.75), y, Inches(1.7), Inches(0.82), k,
                 size=20, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(3.55), y, Inches(9.2), Inches(0.82), v,
                 size=16, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    add_rect(s, Inches(0.45), Inches(5.15), Inches(12.4), Inches(1.45), fill=CLOUD)
    add_text(s, Inches(0.7), Inches(5.3), Inches(12.0), Inches(1.2),
             "启动金 15 万，覆盖头 90 天。之后月度 8 万。散户首年租金 8%，整层 / 大客户 12%。\n"
             "启动金成交后抵佣金。常规中介约年租金 16%，我方大约一半。",
             size=15, color=INK)

    total = len(SLIDES)
    for sl in SLIDES:
        for shape in sl.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text and run.text.strip().endswith(" / 0"):
                        run.text = run.text.replace(" / 0", f" / {total}")

    out = Path(__file__).resolve().parent.parent / "docs" / "advisory" / \
        "deck" / "冠松01楼-政策路径-董事会版.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"✓ Deck written: {out}  ({total} slides)")
    return out


if __name__ == "__main__":
    build()
