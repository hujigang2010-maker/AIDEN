"""
GS · iDrive Hub 招商方案 PPT 生成脚本（商务汇报版 · 16:9 · 35 页）

- 中文字体：WenQuanYi Micro Hei（系统已安装；可改 PingFang/Source Han）
- 主色板：
    深海蓝 #0F2D52 主基调
    智驾蓝 #1F6FEB 强调 / 数据
    金 砂 #C9A24A 重点 / 数字高亮
    云 灰 #F4F6FA 正文背景
    炭 黑 #1B1F2A 正文字
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pathlib import Path

NAVY = RGBColor(0x0F, 0x2D, 0x52)
BLUE = RGBColor(0x1F, 0x6F, 0xEB)
GOLD = RGBColor(0xC9, 0xA2, 0x4A)
CLOUD = RGBColor(0xF4, 0xF6, 0xFA)
INK = RGBColor(0x1B, 0x1F, 0x2A)
GREY = RGBColor(0x6B, 0x73, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD8, 0xDE, 0xE9)
GREEN = RGBColor(0x2F, 0xA3, 0x6F)
RED = RGBColor(0xD0, 0x4A, 0x4A)
PURPLE = RGBColor(0x6B, 0x4A, 0xC0)

CN_FONT = "WenQuanYi Micro Hei"
EN_FONT = "Calibri"

SLIDES = []  # 全局，便于回填总页数


# --------------------------------------------------------------------------
# 工具函数
# --------------------------------------------------------------------------
def set_run(run, text, *, size=14, bold=False, color=INK, italic=False,
            font_cn=CN_FONT, font_en=EN_FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_en
    rPr = run._r.get_or_add_rPr()
    for tag in ("ea", "cs"):
        existing = rPr.find(qn(f"a:{tag}"))
        if existing is not None:
            rPr.remove(existing)
        ea = rPr.makeelement(qn(f"a:{tag}"), {"typeface": font_cn})
        rPr.append(ea)


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, fill=None, line=None,
             italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if fill is not None:
        box.fill.solid(); box.fill.fore_color.rgb = fill
    else:
        box.fill.background()
    if line is not None:
        box.line.color.rgb = line; box.line.width = Pt(0.5)
    else:
        box.line.fill.background()
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        set_run(run, ln, size=size, bold=bold, color=color, italic=italic)
    return box


def add_rect(slide, x, y, w, h, *, fill=NAVY, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(0.5)
    shp.shadow.inherit = False
    shp.text_frame.text = ""
    return shp


def add_round(slide, x, y, w, h, text="", *, fill=BLUE, color=WHITE,
              size=12, bold=True, align=PP_ALIGN.CENTER):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = align
    set_run(p.add_run(), text, size=size, bold=bold, color=color)
    return shp


def add_table(slide, x, y, w, h, header, rows, *, header_fill=NAVY,
              header_color=WHITE, zebra=(WHITE, CLOUD), header_size=11,
              body_size=10, col_widths=None, body_align=PP_ALIGN.LEFT):
    cols = len(header); n_rows = len(rows) + 1
    ts = slide.shapes.add_table(n_rows, cols, x, y, w, h)
    table = ts.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw
    for j, htxt in enumerate(header):
        cell = table.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Pt(4); tf.margin_right = Pt(4)
        tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
        tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        set_run(p.add_run(), htxt, size=header_size, bold=True, color=header_color)
    for i, row in enumerate(rows, start=1):
        bg = zebra[(i - 1) % 2]
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Pt(4); tf.margin_right = Pt(4)
            tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
            tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = body_align
            set_run(p.add_run(), str(val), size=body_size, color=INK)
    return ts


def add_chrome(slide, prs, *, page_no, phase_label="", page_title="", subtitle=""):
    sw, sh = prs.slide_width, prs.slide_height
    add_rect(slide, 0, 0, sw, Emu(380000), fill=NAVY)
    if phase_label:
        add_round(slide, Inches(0.5), Inches(0.18), Inches(2.0), Inches(0.36),
                  phase_label, fill=GOLD, color=NAVY, size=11, bold=True)
    if page_title:
        add_text(slide, Inches(2.7), Inches(0.10), Inches(8.5), Inches(0.55),
                 page_title, size=22, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(2.7), Inches(0.55), Inches(8.5), Inches(0.30),
                 subtitle, size=11, color=CLOUD, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 0, sh - Emu(80000), sw, Emu(80000), fill=GOLD)
    add_text(slide, Inches(0.5), Inches(7.05), Inches(8.0), Inches(0.30),
             "GS · iDrive Hub · 冠松静安智能驾驶研发中心 · 招商方案 v1.0",
             size=9, color=GREY)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(2.0), Inches(0.30),
             f"{page_no} / 0", size=9, color=GREY, align=PP_ALIGN.RIGHT)


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    SLIDES.append(s)
    return s


# --------------------------------------------------------------------------
# 复用页：章节扉页
# --------------------------------------------------------------------------
def section_cover(prs, idx_label, title_cn, title_en, points):
    s = new_slide(prs)
    sw, sh = prs.slide_width, prs.slide_height
    add_rect(s, 0, 0, sw, sh, fill=NAVY)
    add_rect(s, 0, Inches(3.7), sw, Emu(40000), fill=GOLD)
    add_text(s, Inches(0.8), Inches(1.0), Inches(8), Inches(0.6),
             idx_label, size=18, bold=True, color=GOLD)
    add_text(s, Inches(0.8), Inches(1.7), Inches(11), Inches(1.3),
             title_cn, size=44, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.0), Inches(11), Inches(0.6),
             title_en, size=18, color=CLOUD, italic=True)
    for i, p in enumerate(points):
        y = Inches(4.2) + Inches(0.55) * i
        add_round(s, Inches(0.8), y, Inches(0.4), Inches(0.4), str(i + 1),
                  fill=GOLD, color=NAVY, size=14, bold=True)
        add_text(s, Inches(1.4), y, Inches(11), Inches(0.4), p,
                 size=15, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    return s


# --------------------------------------------------------------------------
# 链主一页纸（统一模板）
# --------------------------------------------------------------------------
def anchor_one_pager(prs, *, idx, label, brand, title, why_anchor, what_lacks,
                     unique_offer, deal_terms, next_steps):
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label=f"Phase 2 · 任务 3 · 提案 {idx}",
               page_title=f"{brand} · 一页纸定制提案",
               subtitle=title)
    # 顶部："为什么选静安、为什么选冠松"两块
    top_y = Inches(1.20); top_h = Inches(1.55)
    add_rect(s, Inches(0.5), top_y, Inches(6.15), top_h, fill=NAVY)
    add_text(s, Inches(0.7), top_y + Inches(0.10), Inches(5.8), Inches(0.4),
             f"为什么是 {label}", size=14, bold=True, color=GOLD)
    add_text(s, Inches(0.7), top_y + Inches(0.50), Inches(5.8), Inches(1.0),
             why_anchor, size=11, color=WHITE)

    add_rect(s, Inches(6.85), top_y, Inches(6.15), top_h, fill=BLUE)
    add_text(s, Inches(7.05), top_y + Inches(0.10), Inches(5.8), Inches(0.4),
             "他正缺什么", size=14, bold=True, color=GOLD)
    add_text(s, Inches(7.05), top_y + Inches(0.50), Inches(5.8), Inches(1.0),
             what_lacks, size=11, color=WHITE)

    # 中部：独家供给（4 列卡片）
    mid_y = Inches(2.95)
    add_text(s, Inches(0.5), mid_y, Inches(12.5), Inches(0.4),
             "静安 + 冠松 · 独家供给", size=14, bold=True, color=NAVY)
    add_rect(s, Inches(0.5), mid_y + Inches(0.40), Inches(12.5), Emu(20000), fill=GOLD)
    cw = Inches(3.10); gx = Inches(0.05); cy = mid_y + Inches(0.55); ch = Inches(1.45)
    colors = [BLUE, NAVY, GOLD, GREEN]
    for i, (t, d) in enumerate(unique_offer):
        x = Inches(0.5) + (cw + gx) * i
        add_rect(s, x, cy, cw, Inches(0.45), fill=colors[i % 4])
        add_text(s, x, cy + Inches(0.05), cw, Inches(0.35), t,
                 size=12, bold=True, color=WHITE if colors[i % 4] != GOLD else NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, x, cy + Inches(0.45), cw, ch - Inches(0.45), fill=CLOUD, line=LINE)
        add_text(s, x + Inches(0.10), cy + Inches(0.55), cw - Inches(0.2), ch - Inches(0.55),
                 d, size=10, color=INK)

    # 报价权益
    deal_y = Inches(4.95)
    add_text(s, Inches(0.5), deal_y, Inches(12.5), Inches(0.4),
             "报价与权益（草案）", size=14, bold=True, color=NAVY)
    add_rect(s, Inches(0.5), deal_y + Inches(0.40), Inches(12.5), Emu(20000), fill=GOLD)
    add_rect(s, Inches(0.5), deal_y + Inches(0.55), Inches(8.0), Inches(1.50),
             fill=CLOUD, line=LINE)
    for i, line in enumerate(deal_terms):
        y = deal_y + Inches(0.65) + Inches(0.22) * i
        add_round(s, Inches(0.65), y + Inches(0.05), Inches(0.18), Inches(0.18), "●",
                  fill=BLUE, color=WHITE, size=8)
        add_text(s, Inches(0.90), y, Inches(7.6), Inches(0.22), line,
                 size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # 下一步
    add_rect(s, Inches(8.65), deal_y + Inches(0.55), Inches(4.35), Inches(1.50),
             fill=NAVY)
    add_text(s, Inches(8.80), deal_y + Inches(0.62), Inches(4.1), Inches(0.32),
             "下一步行动", size=12, bold=True, color=GOLD)
    for i, ln in enumerate(next_steps):
        y = deal_y + Inches(0.95) + Inches(0.32) * i
        add_round(s, Inches(8.80), y + Inches(0.04), Inches(0.7), Inches(0.22),
                  ln[0], fill=GOLD, color=NAVY, size=10, bold=True)
        add_text(s, Inches(9.55), y, Inches(3.4), Inches(0.30), ln[1],
                 size=10, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    return s


# --------------------------------------------------------------------------
# 主流程
# --------------------------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    sw, sh = prs.slide_width, prs.slide_height

    # ============ 1. 封面 ============
    s = new_slide(prs)
    add_rect(s, 0, 0, sw, sh, fill=NAVY)
    add_rect(s, 0, Inches(5.5), sw, Inches(0.06), fill=GOLD)
    deco = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(-2.0), Inches(8.0), Inches(8.0))
    deco.fill.solid(); deco.fill.fore_color.rgb = BLUE
    deco.line.fill.background(); deco.shadow.inherit = False
    add_text(s, Inches(0.8), Inches(0.6), Inches(6), Inches(0.5),
             "GS · iDrive Hub", size=18, bold=True, color=GOLD)
    add_text(s, Inches(0.8), Inches(1.6), Inches(11), Inches(1.6),
             "01# 新建研发楼\n智能驾驶研发与总部楼", size=44, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.8), Inches(11), Inches(0.6),
             "让中心城区跑通智能驾驶最后一公里",
             size=22, color=CLOUD, italic=True)
    add_text(s, Inches(0.8), Inches(4.6), Inches(11), Inches(0.5),
             "永和社区 075b-07 地块 · 1.5 万㎡ · 9F · 高 44.95 m · C6 教育科研用地 · 绿建二星",
             size=13, color=CLOUD)
    add_round(s, Inches(0.8), Inches(6.2), Inches(2.4), Inches(0.45),
              "v1.1 · 基于设计 PDF 实测", fill=GOLD, color=NAVY, size=12, bold=True)
    add_text(s, Inches(3.4), Inches(6.2), Inches(8), Inches(0.45),
             "汇报对象：集团董事会 / 静安区政府 / 链主企业",
             size=11, color=CLOUD, anchor=MSO_ANCHOR.MIDDLE)

    # ============ 2. 议程 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="议程",
               page_title="汇报议程", subtitle="约 35 分钟形成完整认知")
    items = [
        ("01", "执行摘要 · 项目实测 · SWOT", "01# 研发楼 · 1.5 万方真实口径"),
        ("02", "Phase 1 · 策略与定位", "产业研究 + 9F 垂直拆分 + 绿色 ESG"),
        ("03", "Phase 2 · 招商执行", "链主 5 张一页纸 + 漏斗 + 政府"),
        ("04", "Phase 3 · 品牌与活动", "9 月发布会 + 年度活动 + 媒体"),
        ("05", "Phase 4 · 商业条款", "四档 + 服务包 + 报价 + 财务"),
        ("06", "Phase 5 · 落地推进", "12 个月甘特 + 团队 + 风险矩阵"),
        ("07", "投决建议", "三个里程碑承诺 + 5 项授权"),
    ]
    y0 = Inches(1.3)
    for i, (no, title, sub) in enumerate(items):
        y = y0 + Inches(0.75) * i
        add_round(s, Inches(0.8), y, Inches(0.7), Inches(0.55), no,
                  fill=NAVY, color=WHITE, size=18, bold=True)
        add_text(s, Inches(1.7), y, Inches(6), Inches(0.55), title,
                 size=18, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(7.7), y, Inches(5), Inches(0.55), sub,
                 size=13, color=GREY, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(0.8), y + Inches(0.62), Inches(11.7), Emu(15000), fill=LINE)

    # ============ 3. 一页摘要 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="01 · 核心结论",
               page_title="一页看懂：项目战略与三年目标",
               subtitle="GS · iDrive Hub · 中心城区智能驾驶研发与总部首选地")
    card_w = Inches(2.95); card_h = Inches(1.55); gap = Inches(0.15)
    cards = [
        ("15,153 ㎡", "地上 9F · 高 44.95 m\n+ 地下 6,993 ㎡", BLUE),
        ("5.4–5.7 m", "3–4F 净层高\n中心城区硬科技稀缺", GOLD),
        ("绿建二星", "装配式 100% · 540 ㎡ 光伏\n碳排放 −48.93%", GREEN),
        ("92%", "Y3 入驻率 · 链主≥1\n入驻企业 12–18 家", NAVY),
    ]
    x = Inches(0.5)
    for i, (big, sub, color) in enumerate(cards):
        bx = x + (card_w + gap) * i
        add_rect(s, bx, Inches(1.1), card_w, card_h, fill=color)
        add_text(s, bx, Inches(1.18), card_w, Inches(0.85), big,
                 size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, bx, Inches(1.95), card_w, Inches(0.65), sub,
                 size=11, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.5), Inches(2.95), Inches(12.3), Inches(0.4),
             "战略定位 · 三圈层客群 · 四大差异化壁垒",
             size=15, bold=True, color=NAVY)
    add_rect(s, Inches(0.5), Inches(3.40), Inches(12.3), Emu(20000), fill=GOLD)
    block_y = Inches(3.55); block_h = Inches(2.6)
    add_rect(s, Inches(0.5), block_y, Inches(4.0), block_h, fill=CLOUD)
    add_text(s, Inches(0.65), block_y + Inches(0.1), Inches(3.7), Inches(0.4),
             "战略定位", size=14, bold=True, color=NAVY)
    add_text(s, Inches(0.65), block_y + Inches(0.55), Inches(3.7), Inches(2.0),
             "L4 城市 NOA + 车路云一体化 + 智驾后市场\n\n研发总部 + 测试验证 + 产业服务\n三位一体园区\n\n品牌主张：让中心城区跑通\n智能驾驶最后一公里", size=11, color=INK)
    add_rect(s, Inches(4.65), block_y, Inches(4.0), block_h, fill=CLOUD)
    add_text(s, Inches(4.80), block_y + Inches(0.1), Inches(3.7), Inches(0.4),
             "垂直三层客群", size=14, bold=True, color=NAVY)
    add_text(s, Inches(4.80), block_y + Inches(0.55), Inches(3.7), Inches(2.0),
             "8–9F · 链主总部\n4.2–4.3 m · 1–2 家\n\n3–4F · 核心研发\n5.4–5.7 m 高层高 · 4–6 家\n\n6–7F · 算法软件\n4.2 m · 6–10 家",
             size=11, color=INK)
    add_rect(s, Inches(8.80), block_y, Inches(4.0), block_h, fill=NAVY)
    add_text(s, Inches(8.95), block_y + Inches(0.1), Inches(3.7), Inches(0.4),
             "五大差异化壁垒", size=14, bold=True, color=GOLD)
    add_text(s, Inches(8.95), block_y + Inches(0.55), Inches(3.7), Inches(2.0),
             "① 3–4F 5.4–5.7 m 硬科技层高\n\n② 绿建二星 + 装配式 + 光伏\n\n③ 静安一企一策政策包\n\n④ 冠松汽车后市场资源闭环\n\n⑤ 三段式测试方案 (轻资产)",
             size=11, color=WHITE)
    add_text(s, Inches(0.5), Inches(6.30), Inches(12.3), Inches(0.4),
             "关键举措：5 链主 · 300 生态 · 5 中介 · 1 场 9 月发布会 · 4 档商业模式",
             size=12, bold=True, color=NAVY)

    # ============ 4. 项目区位 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="01 · 项目区位",
               page_title="项目区位 · 中心城区独栋园区的稀缺性",
               subtitle="地铁双线直达 · 30 min 触达虹桥/浦东 · 1.5 km 环内可申请测试延伸路段")
    # 左：放射状区位图
    add_text(s, Inches(0.5), Inches(1.15), Inches(7.5), Inches(0.4),
             "区位辐射示意（GS · iDrive Hub 为中心）", size=14, bold=True, color=NAVY)
    map_x, map_y, map_w, map_h = Inches(0.5), Inches(1.55), Inches(7.5), Inches(5.3)
    add_rect(s, map_x, map_y, map_w, map_h, fill=CLOUD, line=LINE)
    cx = map_x + map_w / 2
    cy = map_y + map_h / 2
    # 中心点
    center = s.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(0.45), cy - Inches(0.45),
                                Inches(0.9), Inches(0.9))
    center.fill.solid(); center.fill.fore_color.rgb = GOLD
    center.line.color.rgb = NAVY
    center.line.width = Pt(2)
    center.shadow.inherit = False
    tf = center.text_frame
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(), "GS\nHub", size=10, bold=True, color=NAVY)
    # 8 个方向地点
    pois = [
        ("南京西路 CBD", -2.6, -1.3, "10 min"),
        ("人民广场", -2.0, -2.1, "12 min"),
        ("陆家嘴", 2.6, -1.3, "25 min"),
        ("张江 / AI 岛", 3.2, 0.8, "40 min"),
        ("虹桥枢纽 / 机场", -3.0, 1.0, "30 min"),
        ("浦东机场", 3.2, 1.8, "55 min"),
        ("嘉定汽车城", -1.0, -2.4, "45 min"),
        ("临港滴水湖", 2.4, 2.2, "75 min"),
    ]
    for name, dx, dy, t in pois:
        px = cx + Inches(dx) - Inches(0.95)
        py = cy + Inches(dy) - Inches(0.32)
        # 连线（用细矩形作箭杆）
        line_shp = s.shapes.add_connector(1, cx, cy, px + Inches(0.95), py + Inches(0.32))
        line_shp.line.color.rgb = BLUE
        line_shp.line.width = Pt(1)
        # 节点圆角
        add_round(s, px, py, Inches(1.9), Inches(0.34), name,
                  fill=NAVY, color=WHITE, size=10, bold=True)
        add_round(s, px + Inches(1.92), py, Inches(0.6), Inches(0.34),
                  t, fill=GOLD, color=NAVY, size=9, bold=True)

    # 右：交通配套数据卡
    add_text(s, Inches(8.4), Inches(1.15), Inches(4.7), Inches(0.4),
             "交通与配套（步行/车程口径）", size=14, bold=True, color=NAVY)
    facts = [
        ("地铁", "双线直达 · 园区主入口 5 min", BLUE),
        ("公交", "12 条公交线 · 多个站点环抱", BLUE),
        ("机场", "虹桥 30 min · 浦东 55 min", NAVY),
        ("高铁", "上海站 15 min · 虹桥站 30 min", NAVY),
        ("CBD", "南京西路 10 min · 人民广场 12 min", GOLD),
        ("产业", "嘉定 45 min · 张江 40 min · 临港 75 min", GOLD),
        ("人才", "周边 5 km 内 12 所高校 / 60 万 IT 白领", GREEN),
        ("住居", "区内 8 个人才公寓社区 + 高端住宅带", GREEN),
    ]
    for i, (t, d, c) in enumerate(facts):
        y = Inches(1.55) + Inches(0.65) * i
        add_rect(s, Inches(8.4), y, Inches(0.18), Inches(0.55), fill=c)
        add_rect(s, Inches(8.6), y, Inches(4.5), Inches(0.55), fill=CLOUD, line=LINE)
        add_text(s, Inches(8.75), y + Inches(0.05), Inches(1.0), Inches(0.45),
                 t, size=12, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(9.80), y + Inches(0.05), Inches(3.25), Inches(0.45),
                 d, size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # ============ 5. SWOT 总览 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="01 · SWOT",
               page_title="SWOT 总览 · 把优势变成差异化报价",
               subtitle="把外部机会与内部优势耦合，形成不可复制的 4 道护城河")
    quad_w = Inches(6.15); quad_h = Inches(2.85)
    quads = [
        ("优势 S", "Strengths", BLUE,
         ["3–4F 净高 5.4–5.7 m 中心城区硬科技稀缺",
          "绿建二星 + 装配式 100% + 540 ㎡ 光伏",
          "冠松 4S / 保险 / 二手车后市场资源",
          "C6 教育科研用地 · 政策定位明确",
          "地铁直达 · 30 min 触达虹桥/浦东"]),
        ("劣势 W", "Weaknesses", GOLD,
         ["单体 1.5 万方（地上）· 出租净 ~8,300㎡",
          "无地块内大型户外封闭测试场",
          "C6 用地限制纯商务 / 制造 / 仓储",
          "路测牌照需逐级申请",
          "中心城区运营成本高于郊区"]),
        ("机会 O", "Opportunities", NAVY,
         ["城市 NOA + L3 法规 进入量产爆发期",
          "智驾后市场（保险/改装/二手车）爆发",
          "央企 / 外资 ADAS 总部外迁需求",
          "政府推动「中心城区智算/智驾」走廊",
          "L3 准入试点带来政策红利窗口"]),
        ("威胁 T", "Threats", RED,
         ["嘉定 / 临港政策护城河深",
          "张江 / 徐汇高端租金天花板压缩定价",
          "智驾行业整合期，腰部企业出清",
          "地缘 / 数据合规对外资客户造成不确定",
          "公共预算紧张，专项补贴趋严"]),
    ]
    positions = [(Inches(0.5), Inches(1.20)), (Inches(6.85), Inches(1.20)),
                 (Inches(0.5), Inches(4.20)), (Inches(6.85), Inches(4.20))]
    for (px, py), (t, en, c, bullets) in zip(positions, quads):
        add_rect(s, px, py, quad_w, quad_h, fill=CLOUD, line=LINE)
        add_rect(s, px, py, quad_w, Inches(0.50), fill=c)
        add_text(s, px + Inches(0.15), py + Inches(0.05), quad_w - Inches(0.3), Inches(0.4),
                 t, size=15, bold=True, color=GOLD if c == NAVY else WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, px + quad_w - Inches(2.0), py + Inches(0.05), Inches(1.8), Inches(0.4),
                 en, size=11, color=WHITE, italic=True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        for i, b in enumerate(bullets):
            yy = py + Inches(0.65) + Inches(0.42) * i
            add_round(s, px + Inches(0.20), yy + Inches(0.08), Inches(0.18), Inches(0.18), "●",
                      fill=c, color=WHITE, size=8)
            add_text(s, px + Inches(0.45), yy, quad_w - Inches(0.6), Inches(0.4), b,
                     size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # ============ 6. Phase 1 扉页 ============
    section_cover(prs, "PHASE 1", "策略与定位",
                  "Strategy & Positioning · 奠基",
                  ["任务 1 · 产业定位研究与竞品对标",
                   "任务 2 · 空间功能规划（A~E 栋 + 户外测试区）"])

    # ============ 7. 产业链图谱 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 1 · 任务 1",
               page_title="智能驾驶产业链六层图谱",
               subtitle="客群匹配的『地图』：从基础设施到应用运营")
    layers = [
        ("L6 应用与运营", "Robotaxi · Robobus · Robotruck · 矿区港口 · 末端配送", BLUE),
        ("L5 整车与方案", "车企 ADAS/AD 部门 · 全栈方案商（华为/Momenta/...）", NAVY),
        ("L4 软件与算法", "感知 · 规控 · 定位 · 端到端大模型 · 仿真 · 数据闭环", BLUE),
        ("L3 域控与计算", "智驾域控 · 中央计算 · 芯片（地平线/黑芝麻/英伟达）", NAVY),
        ("L2 传感器", "激光雷达 · 4D 毫米波 · 摄像头 · 高精地图 · IMU/RTK", BLUE),
        ("L1 基础设施", "路侧 RSU · V2X · 5G/卫星通信 · 高精地图测绘资质", NAVY),
    ]
    bar_x = Inches(0.7); bar_w = Inches(8.5); bar_h = Inches(0.65)
    for i, (lbl, desc, c) in enumerate(layers):
        y = Inches(1.25) + (bar_h + Inches(0.10)) * i
        add_rect(s, bar_x, y, Inches(1.6), bar_h, fill=GOLD)
        add_text(s, bar_x, y, Inches(1.6), bar_h, lbl,
                 size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, bar_x + Inches(1.6), y, bar_w - Inches(1.6), bar_h, fill=c)
        add_text(s, bar_x + Inches(1.7), y, bar_w - Inches(1.7), bar_h, desc,
                 size=11, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    side_x = Inches(9.5); side_w = Inches(3.4)
    side = [("资本 / 孵化", "国调/中金/高瓴/园区基金"),
            ("测试 / 认证", "中汽研 · TüV · CATARC"),
            ("后市场服务", "保险 · 维修 · 改装 · 二手车 (冠松)")]
    for i, (t, d) in enumerate(side):
        y = Inches(1.6) + Inches(1.45) * i
        add_rect(s, side_x, y, side_w, Inches(1.25), fill=CLOUD, line=LINE)
        add_text(s, side_x + Inches(0.15), y + Inches(0.1), side_w - Inches(0.3), Inches(0.4),
                 t, size=14, bold=True, color=NAVY)
        add_text(s, side_x + Inches(0.15), y + Inches(0.5), side_w - Inches(0.3), Inches(0.7),
                 d, size=11, color=INK)
    add_text(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.4),
             "→ A 栋承接 L5 链主 · B 栋承接 L3/L4 · C 栋承接仿真/数据 · D 栋承接服务集群 · E 栋承接冠松后市场",
             size=12, bold=True, color=BLUE)

    # ============ 8. 竞品对标 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 1 · 任务 1",
               page_title="静安区内 + 跨区竞品对标",
               subtitle="核心结论：中心城区独栋+测试区组合稀缺，租金锚定 6.5–8.5 元/㎡·天")
    header = ["园区", "区位", "主导产业", "租金 (元/㎡·天)", "入驻率", "政策亮点", "对我方"]
    rows = [
        ["市北高新 · 大数据基地", "静安", "大数据/AI/智算", "6.0–8.5", "~88%", "区级 5%~10% 租补 · 公寓 800 套", "无智驾产业链 / 无测试区"],
        ["静安国际中心片区", "静安", "总部金融", "9.0–13.0", "~90%", "总部经济奖励", "无产业属性 / 不可改造"],
        ["大宁国际办公", "静安", "商务/文创", "6.5–8.0", "~85%", "一般", "非产业园定位"],
        ["嘉定 · 创新中心 (新能港)", "嘉定", "智驾全产业链", "3.5–5.0", "~95%", "国家级测试场 + 牌照", "牌照不可复制 · 我方区位优"],
        ["临港 · AI 创新港", "浦东", "L4/Robotaxi", "3.0–4.5", "~80%", "装补 800–1500/㎡ · 全域开放", "政策强 · 距市区 1.5h"],
        ["张江 · IC 设计园 / AI 岛", "浦东", "芯片/AI", "6.5–8.5", "~92%", "张江专项 + IC 补贴", "算力强 · 智驾整车弱"],
        ["徐汇 · 西岸 AI 走廊", "徐汇", "AI 大模型", "8.5–11.0", "~95%", "模速空间补贴", "单价高 · 无独栋"],
    ]
    add_table(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(4.6), header, rows,
              col_widths=[Inches(2.4), Inches(0.8), Inches(1.5), Inches(1.4),
                          Inches(0.9), Inches(2.7), Inches(2.6)])
    add_rect(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.95), fill=NAVY)
    add_text(s, Inches(0.7), Inches(6.05), Inches(12.0), Inches(0.4),
             "我方差异化定价策略", size=13, bold=True, color=GOLD)
    add_text(s, Inches(0.7), Inches(6.40), Inches(12.0), Inches(0.55),
             "基础租金（6.5–8.5）+ 政策返还（区级留成 80% 三年返）+ 服务积分（测试/算力/招聘） — 等效净价低于南西核心 25%",
             size=11, color=WHITE)

    # ============ 9. 单楼垂直功能拆分 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 1 · 任务 2",
               page_title="01# 研发楼 · 9F 垂直功能拆分",
               subtitle="链主总部(8–9F) · 算法软件(6–7F) · 公共配套(5F) · 核心研发(3–4F) · 大堂展厅(1–2F)")
    base_x = Inches(0.5); base_y = Inches(1.2)
    plot_w = Inches(6.5); plot_h = Inches(5.4)
    add_rect(s, base_x, base_y, plot_w, plot_h, fill=CLOUD, line=LINE)
    add_text(s, base_x, base_y + Inches(0.05), plot_w, Inches(0.3),
             "01# 研发楼 · 垂直叠加（北侧立面示意）", size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # 楼层数据：净层高 (m), 标签, 描述, 颜色
    floors_v = [
        ("RF", "—",   "屋顶花园 + 540 ㎡ 光伏", GREEN),
        ("9F", "4.3", "总部 · CEO/战略/投资 (链主)", GOLD),
        ("8F", "4.2", "总部 · 产品/商务/GR (链主)", GOLD),
        ("7F", "4.2", "标准研发 · 算法/软件", BLUE),
        ("6F", "4.2", "标准研发 · 算法/数据", BLUE),
        ("5F", "4.2", "公共配套 · 餐厅 + 健身 (不出租)", GREY),
        ("4F", "5.4", "核心研发 · 高层高硬件/座舱", PURPLE),
        ("3F", "5.7", "核心研发 · 联合实验室", PURPLE),
        ("2F", "6.3", "展厅 + 食堂 (1F+2F 可冠名)", NAVY),
        ("1F", "5.7", "大堂 + 展厅 + 接待 (品牌门面)", NAVY),
        ("B1", "—",   "物业 + 设备", LINE),
        ("B2", "—",   "停车 108 位 (智驾专车 30–50)", LINE),
    ]
    fly = base_y + Inches(0.45); flh = Inches(0.39); flx = base_x + Inches(0.15)
    flw = plot_w - Inches(0.3)
    for i, (lvl, h, desc, c) in enumerate(floors_v):
        y = fly + (flh + Inches(0.02)) * i
        add_rect(s, flx, y, Inches(0.55), flh, fill=c)
        add_text(s, flx, y, Inches(0.55), flh, lvl,
                 size=10, bold=True, color=NAVY if c in (GOLD, LINE) else WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, flx + Inches(0.6), y, Inches(0.65), flh, fill=CLOUD, line=LINE)
        add_text(s, flx + Inches(0.6), y, Inches(0.65), flh, h,
                 size=10, bold=True, color=GREY if h == "—" else NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, flx + Inches(1.30), y, flw - Inches(1.30), flh,
                 fill=WHITE, line=LINE)
        add_text(s, flx + Inches(1.4), y, flw - Inches(1.45), flh, desc,
                 size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, base_x, base_y + plot_h - Inches(0.30), plot_w, Inches(0.25),
             "层 / 净层高 (m) / 功能",
             size=9, color=GREY, align=PP_ALIGN.CENTER, italic=True)

    # 右侧 · 真实指标表
    header = ["关键指标", "数值"]
    rows = [
        ["地上建筑面积", "15,152.75 ㎡（无生产）"],
        ["地下建筑面积", "6,992.87 ㎡"],
        ["建筑高度", "44.95 m"],
        ["楼层", "地上 9F · 地下 2F"],
        ["可出租净面积（地上）", "约 8,300 ㎡"],
        ["1F+2F 冠名空间", "约 1,100 ㎡"],
        ["地下停车", "108 个（3 普 + 105 机械）"],
        ["智驾专车规划车位", "30–50 个（含充电）"],
        ["用地性质", "C6 教育科研设计用地"],
        ["建筑高度内荷载", "标准 ≥ 4.0 kN/㎡"],
        ["建议加固（3–4F）", "≥ 7.5 kN/㎡"],
        ["供电预留", "标准 ≥ 200 W/㎡，3–4F ≥ 350"],
    ]
    add_table(s, Inches(7.3), Inches(1.2), Inches(5.5), Inches(5.4),
              header, rows,
              col_widths=[Inches(2.6), Inches(2.9)],
              header_size=11, body_size=10)

    # ============ 10. 8–9F 链主总部深度 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 1 · 任务 2",
               page_title="8–9F · 链主总部楼层（『顶部冠名』）",
               subtitle="2 整层共 ~3,366 ㎡ · 净高 4.2–4.3 m · 起始 5.5–6.5 元/㎡·天 · 6 年长租")
    # 左：8/9F 平面分区
    add_text(s, Inches(0.5), Inches(1.15), Inches(5.5), Inches(0.4),
             "9F · CEO 与战略层（约 1,683 ㎡）", size=13, bold=True, color=NAVY)
    add_rect(s, Inches(0.5), Inches(1.55), Inches(5.5), Inches(2.0), fill=CLOUD, line=LINE)
    parts9 = [
        ("CEO 室+露台",  Inches(0.6),  Inches(1.65), Inches(1.6), Inches(1.0), GOLD, "16%"),
        ("战略 + 投资", Inches(2.3),  Inches(1.65), Inches(1.6), Inches(1.0), BLUE, "25%"),
        ("高管接待", Inches(4.0),   Inches(1.65), Inches(1.9), Inches(1.0), NAVY, "18%"),
        ("多功能厅", Inches(0.6),   Inches(2.70), Inches(2.5), Inches(0.85), PURPLE, "22%"),
        ("开放办公 + 会议", Inches(3.2), Inches(2.70), Inches(2.7), Inches(0.85), BLUE, "19%"),
    ]
    for txt, x, y, w, h, c, p in parts9:
        add_rect(s, x, y, w, h, fill=c)
        add_text(s, x, y, w, h, txt + "\n" + p,
                 size=9, bold=True, color=WHITE if c != GOLD else NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.5), Inches(3.7), Inches(5.5), Inches(0.4),
             "8F · 产品 / 商务 / GR 层（约 1,683 ㎡）", size=13, bold=True, color=NAVY)
    add_rect(s, Inches(0.5), Inches(4.10), Inches(5.5), Inches(2.0), fill=CLOUD, line=LINE)
    parts8 = [
        ("开放办公区 (产品+研发管理)", Inches(0.6), Inches(4.20), Inches(3.4), Inches(1.0), BLUE, "45%"),
        ("会议室集群", Inches(4.1), Inches(4.20), Inches(1.8), Inches(1.0), NAVY, "20%"),
        ("商务 + GR 独立办公",  Inches(0.6), Inches(5.25), Inches(2.5), Inches(0.80), GOLD, "22%"),
        ("茶水/休闲/打印",  Inches(3.2), Inches(5.25), Inches(2.7), Inches(0.80), GREEN, "13%"),
    ]
    for txt, x, y, w, h, c, p in parts8:
        add_rect(s, x, y, w, h, fill=c)
        add_text(s, x, y, w, h, txt + "\n" + p,
                 size=9, bold=True, color=WHITE if c != GOLD else NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 右：核心卖点 + 签约组合
    add_text(s, Inches(7.0), Inches(1.15), Inches(6), Inches(0.4),
             "核心卖点", size=14, bold=True, color=NAVY)
    sells = [
        ("楼宇冠名权",   "「[链主] · iDrive Tower 01」 5 年", GOLD),
        ("1F 大堂背景墙", "5.7 m 挑空 + 主形象墙 + LED",   BLUE),
        ("9F 屋顶花园", "ESG/媒体大片专属取景日 ≥ 2/年",  NAVY),
        ("地下专属车位", "8–12 个 · 智驾测试车连通 1F",   BLUE),
        ("装配式 100%",  "整层无柱大空间 · 净高 4.2–4.3 m", GREEN),
        ("供电/弱电",    "≥ 200 W/㎡ · POL + 5G-A + Wi-Fi 6E", PURPLE),
    ]
    for i, (t, d, c) in enumerate(sells):
        y = Inches(1.55) + Inches(0.55) * i
        add_rect(s, Inches(7.0), y, Inches(0.18), Inches(0.45), fill=c)
        add_rect(s, Inches(7.2), y, Inches(5.7), Inches(0.45), fill=CLOUD, line=LINE)
        add_text(s, Inches(7.35), y + Inches(0.02), Inches(2.3), Inches(0.41),
                 t, size=11, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(9.65), y + Inches(0.02), Inches(3.2), Inches(0.41),
                 d, size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    add_rect(s, Inches(7.0), Inches(4.95), Inches(5.9), Inches(1.95), fill=NAVY)
    add_text(s, Inches(7.15), Inches(5.0), Inches(5.6), Inches(0.4),
             "签约组合（草案）", size=13, bold=True, color=GOLD)
    deal = ("起始租金 5.5–6.5 元/㎡·天 · 6 年（含免租）\n"
            "免租 12–18 个月 · 装补 800–1,200 元/㎡ (封顶 400 万)\n"
            "区级税收留成 80% 三年返 / 50% 后两年返\n"
            "嘉定/临港封闭测试场会员代办 · 1.5 km 路测延伸段联合申请\n"
            "人才公寓 80–120 套 + 高级研发落户 20 个/年\n"
            "9 月发布会主旨演讲 + 媒体首发 + 楼宇冠名 5 年")
    add_text(s, Inches(7.15), Inches(5.35), Inches(5.6), Inches(1.55),
             deal, size=10, color=WHITE)

    # ============ 11. 3–4F 核心研发 + 6–7F 算法软件 + 1–2F 大堂展厅 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 1 · 任务 2",
               page_title="3–7F · 核心研发 + 算法软件 / 1–2F · 大堂展厅",
               subtitle="3–4F 高层高 5.4–5.7 m 硬科技稀缺空间 · 6–7F 算法标准研发 · 1–2F 品牌门面")
    cards = [
        ("3–4F · 核心研发\n（高层高）", "约 3,366 ㎡ · 5.4–5.7 m · 7.0–8.0 元/㎡·天",
         ["净层高 5.4–5.7 m，中心城区罕见",
          "适配域控 / 座舱 / 传感器集成调试台架",
          "建议加固楼面荷载 ≥ 7.5 kN/㎡",
          "3F 联合实验室共建（链主或 Tier1 命名权）",
          "签约 3+3 年 · 免租 6–10 月 · 装补 500–800/㎡"], PURPLE),
        ("6–7F · 算法软件", "约 3,366 ㎡ · 4.2 m · 7.5–8.5 元/㎡·天",
         ["软件 / 算法 / 数据团队标准研发",
          "适配 L4 算法 / 仿真 / 数据闭环企业",
          "全光网到桌面 (POL) + 5G-A + Wi-Fi 6E",
          "可分割 300–800 ㎡ 单元，6–10 家共存",
          "签约 3+3+1 · 免租 4–8 月 · 装补 300–500/㎡"], BLUE),
        ("1–2F · 大堂 + 展厅", "约 3,371 ㎡ · 5.7 m / 6.3 m · 自留+冠名",
         ["1F 5.7 m 挑空大堂 + 展厅 + 接待 + 咖啡",
          "2F 6.3 m 展厅 + 食堂（兼路演 200 人）",
          "可分拆冠名：大堂主背景墙 / 展厅 / LED 屏",
          "用作智驾产业展厅 + 9 月发布会主场",
          "丁档冠名 5 年 + 政府/媒体路演阵地"], NAVY),
    ]
    cw = Inches(4.10); cy = Inches(1.20); ch = Inches(5.5); gx = Inches(0.10)
    for i, (t, sub, bullets, c) in enumerate(cards):
        x = Inches(0.5) + (cw + gx) * i
        add_rect(s, x, cy, cw, ch, fill=CLOUD, line=LINE)
        add_rect(s, x, cy, cw, Inches(1.05), fill=c)
        add_text(s, x, cy + Inches(0.10), cw, Inches(0.5), t,
                 size=18, bold=True, color=GOLD if c == NAVY else WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, cy + Inches(0.55), cw, Inches(0.4), sub,
                 size=12, color=WHITE, align=PP_ALIGN.CENTER)
        for j, b in enumerate(bullets):
            yy = cy + Inches(1.30) + Inches(0.85) * j
            add_round(s, x + Inches(0.20), yy + Inches(0.18), Inches(0.30), Inches(0.30),
                      str(j + 1), fill=c, color=WHITE if c != GOLD else NAVY,
                      size=12, bold=True)
            add_text(s, x + Inches(0.55), yy, cw - Inches(0.7), Inches(0.85),
                     b, size=11, color=INK)

    # ============ 12. 绿色低碳成绩单 + 三段式测试方案 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 1 · 任务 2",
               page_title="绿色低碳成绩单 + 三段式测试方案",
               subtitle="ESG 卖点 + 政府双碳考核 / 取代自建测试场，『轻资产』高效率")
    # 左：绿色低碳成绩单
    add_text(s, Inches(0.5), Inches(1.15), Inches(6), Inches(0.4),
             "绿色低碳与 ESG 成绩单（PDF 实测）", size=14, bold=True, color=NAVY)
    esg = [
        ("绿建二星",      "上海中心城区智驾园区中等级最高之一", GREEN),
        ("装配式 100%",   "预制率 ≥ 40% · 装配整体式框架核心筒", BLUE),
        ("屋面光伏",     "540 ㎡ · 自发自用 + 预留扩展", GOLD),
        ("透水铺装",     "400 ㎡（海绵城市）", BLUE),
        ("下凹式绿地",   "50 ㎡（海绵城市）", GREEN),
        ("雨水蓄水池",   "75 m³ · 雨水回用", BLUE),
        ("年径流控制率", "66.8% · 高于上海市基础要求", PURPLE),
        ("碳排放强度",   "较 2016 节能基准 −48.93%", RED),
        ("单位降碳",     "−7.06 kgCO₂/(㎡·a)", RED),
    ]
    for i, (t, d, c) in enumerate(esg):
        y = Inches(1.55) + Inches(0.50) * i
        add_rect(s, Inches(0.5), y, Inches(0.18), Inches(0.42), fill=c)
        add_rect(s, Inches(0.7), y, Inches(5.6), Inches(0.42), fill=CLOUD, line=LINE)
        add_text(s, Inches(0.85), y + Inches(0.02), Inches(2.0), Inches(0.38),
                 t, size=11, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(2.85), y + Inches(0.02), Inches(3.4), Inches(0.38),
                 d, size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(0.5), Inches(6.10), Inches(5.8), Inches(0.65), fill=NAVY)
    add_text(s, Inches(0.65), Inches(6.13), Inches(5.6), Inches(0.6),
             "→ 整套指标可直接用于「双碳考核 / 链主 ESG 报告 / 政府高质量发展」三类汇报口径",
             size=10, bold=True, color=GOLD, anchor=MSO_ANCHOR.MIDDLE)

    # 右：三段式测试方案
    add_text(s, Inches(6.65), Inches(1.15), Inches(6.5), Inches(0.4),
             "三段式测试方案（替代自建 1.2 万方测试场）", size=14, bold=True, color=NAVY)
    seg = [
        ("段 1 · 园区内 (静态 + 数据)",
         "3–4F 高层高研发 / 联合实验室\nB2 智驾测试车专用车位 30–50 个\n仿真平台 + 共享算力券 + 数据合规沙盒",
         BLUE, "1F"),
        ("段 2 · 静安区 (1.5 km 路测延伸)",
         "区交警支队 + 市公安 + 区经委协同\n申请「iDrive 静安路测专项」夜间/低峰段\n首批 1.5 km，验证后逐步扩展",
         GOLD, "2F"),
        ("段 3 · 嘉定 / 临港 (会员通道)",
         "与上海智能网联汽车创新中心(嘉定)\n临港 AI 创新港 签战略合作\n链主免费会员 / 生态会员 8–9 折",
         GREEN, "3F"),
    ]
    sy = Inches(1.55); sh2 = Inches(1.55)
    for i, (t, d, c, _) in enumerate(seg):
        y = sy + (sh2 + Inches(0.10)) * i
        add_rect(s, Inches(6.65), y, Inches(6.5), sh2, fill=CLOUD, line=LINE)
        add_rect(s, Inches(6.65), y, Inches(0.18), sh2, fill=c)
        add_round(s, Inches(6.85), y + Inches(0.20), Inches(0.7), Inches(0.40),
                  f"段{i+1}", fill=c, color=WHITE if c != GOLD else NAVY,
                  size=12, bold=True)
        add_text(s, Inches(7.65), y + Inches(0.10), Inches(5.4), Inches(0.45),
                 t, size=12, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(7.65), y + Inches(0.55), Inches(5.4), Inches(0.95),
                 d, size=10, color=INK)

    add_rect(s, Inches(6.65), Inches(6.20), Inches(6.5), Inches(0.55), fill=NAVY)
    add_text(s, Inches(6.80), Inches(6.25), Inches(6.3), Inches(0.45),
             "→ 轻资产 · 不重投自建测试场 · 政府关系撬动 + 战略合作组合",
             size=11, bold=True, color=GOLD, anchor=MSO_ANCHOR.MIDDLE)

    # ============ 13. Phase 2 扉页 ============
    section_cover(prs, "PHASE 2", "招商执行（核心）",
                  "Leasing Execution · Anchor + Ecosystem + Government",
                  ["任务 3 · 链主 TOP5 攻坚 + 5 张一页纸",
                   "任务 4 · 300 家生态招商漏斗",
                   "任务 5 · 政府关系与政策包"])

    # ============ 14. 链主 TOP5 攻坚作战图 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 2 · 任务 3",
               page_title="链主 TOP5 攻坚作战图",
               subtitle="项目总监 + GR 总监 双人出动 · 周一例会 / 季度董事会复盘")
    header = ["#", "链主", "落位", "决策人 / 关键人", "当前阶段", "下一动作", "截止", "状态"]
    rows = [
        ["1", "华为车 BU", "8F+9F 整层 ~3,366㎡", "上海中心负责人 / GTS 上海", "T0 名片建联", "区委书记+市经信委约见", "T+30d", "推进中"],
        ["2", "百度 Apollo", "9F 整层 ~1,683㎡", "IDG 总裁 / 华东负责人", "已发邀约函", "联合投促办赴京拜访", "T+45d", "顺利"],
        ["3", "小鹏汽车", "6F+7F 算法 ~3,366㎡", "李力耘 / 上海负责人", "公司层初接触", "邀请出席 9 月发布会", "T+60d", "推进中"],
        ["4", "地平线", "3F 核心研发 ~1,683㎡", "余凯 / 黄畅 / 政企", "高层已建联", "邀请参观高层高研发层", "T+30d", "顺利"],
        ["5", "Momenta", "8F 整层 ~1,683㎡", "曹旭东 / GR VP", "已建联，等会面", "曹 CEO + 静安区长会面", "T+45d", "推进中"],
    ]
    add_table(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.6),
              header, rows,
              col_widths=[Inches(0.4), Inches(1.5), Inches(1.4), Inches(2.0),
                          Inches(1.6), Inches(2.6), Inches(1.0), Inches(1.8)])
    add_text(s, Inches(0.5), Inches(4.0), Inches(12), Inches(0.4),
             "每家「一页纸」定制提案 · 四要点结构（详见接下来 5 页）",
             size=14, bold=True, color=NAVY)
    add_rect(s, Inches(0.5), Inches(4.42), Inches(12.3), Emu(20000), fill=GOLD)
    quad_y = Inches(4.6); quad_h2 = Inches(2.3)
    quads2 = [
        ("决策人 / 关键人", "30 秒电梯演讲 → CEO/总裁层；GR 把区领导接口排进议程", BLUE),
        ("正缺什么", "中心城区落点 · 测试场地 · 后市场闭环 · 政企协同", NAVY),
        ("独家供给", "A 栋整栋冠名 · 户外测试区 · 静安一企一策 · 冠松 4S 流量", GOLD),
        ("报价权益", "起始租 4.5–5.5 · 免租 18–24m · 装补 1,000–1,500/㎡ · 留成 80% 三年返", GREEN),
    ]
    qw = Inches(2.95); gx = Inches(0.15)
    for i, (t, d, c) in enumerate(quads2):
        x = Inches(0.5) + (qw + gx) * i
        add_rect(s, x, quad_y, qw, quad_h2, fill=c)
        add_text(s, x + Inches(0.15), quad_y + Inches(0.15), qw - Inches(0.3), Inches(0.5),
                 t, size=14, bold=True, color=GOLD if c == NAVY else WHITE)
        add_text(s, x + Inches(0.15), quad_y + Inches(0.7), qw - Inches(0.3), quad_h2 - Inches(0.85),
                 d, size=11, color=WHITE)

    # ============ 15-19. 5 张链主一页纸 ============
    anchor_one_pager(prs, idx=1, label="华为", brand="华为车 BU",
        title="静安智能驾驶上海第二中心 — 城市 NOA 与车路云一体化的中心城区样板",
        why_anchor="鸿蒙智行 + 引望 = 中国智驾最具号召力的链主。\n上海中心已落地金桥，需要中心城区「展示 + 客户接待 + 政府协同」的第二极。",
        what_lacks="中心城区客户接待与政商展示\n城市 NOA 路测延伸路段\n鸿蒙生态后市场闭环\n政企协同的高效接口",
        unique_offer=[
            ("A 栋整栋冠名", "「华为 · iDrive Tower」南京西路 30min 直达"),
            ("户外测试 + 路测延伸", "1.5 km 测试延伸路段 联合公安申请"),
            ("冠松 4S 后市场", "鸿蒙智行华东最大经销/售后协同伙伴"),
            ("静安一企一策", "区委书记直接对接 · 一窗通办"),
        ],
        deal_terms=[
            "8F + 9F 整层 ~3,366 ㎡ · 签约 6 年",
            "起始租金 5.8 元/㎡·天 · 3 年一调（CPI+0.5%）",
            "免租 15 个月 · 装补 1,000 元/㎡（封顶 336 万）",
            "区级税收留成 80% 三年返 · 1F 大堂冠名 5 年",
            "嘉定/临港测试场免费会员 · 人才公寓 100 套",
        ],
        next_steps=[
            ("T+15d", "见余承东、王军，约上海中心负责人"),
            ("T+30d", "上海中心负责人 + 静安区委书记互访"),
            ("T+90d", "意向条款书 (Term Sheet) 签署"),
        ])

    anchor_one_pager(prs, idx=2, label="百度", brand="百度 Apollo",
        title="上海中心城区「萝卜快跑 + 城市 NOA」双总部基地",
        why_anchor="Apollo 是中国 L4 Robotaxi 商业化最深的玩家；萝卜快跑已多城运营。\n上海中心城区是 Apollo 尚未补齐的关键拼图。",
        what_lacks="上海中心城区 Robotaxi 运营基地\n城市 NOA 数据闭环\n政府关系与牌照绿通\n商务展示与媒体首发",
        unique_offer=[
            ("户外测试区", "1.5 km 延伸路段 + 冠松车队夜间充电"),
            ("数据合规沙盒", "冠松后市场数据池 + 上海数交所沙盒"),
            ("Apollo 专项工作组", "区交警 + 经委 + 投促办承诺成立"),
            ("A 栋路演中心", "1F 600 ㎡ + 冠名首发权"),
        ],
        deal_terms=[
            "9F 整层 ~1,683 ㎡ · 签约 5 年（可叠 8F 至 ~3,366 ㎡）",
            "起始租金 5.5 元/㎡·天",
            "免租 12 个月 · 装补 1,000 元/㎡",
            "1F 大堂冠名 + 与冠松共建「Apollo · GS 出行实验室」",
            "牌照绿色通道：区交警 + 市经信委协作",
        ],
        next_steps=[
            ("T+10d", "联合静安投促办赴京拜会 IDG"),
            ("T+45d", "王云鹏（或代表）静安考察"),
            ("T+120d", "Term Sheet 签署"),
        ])

    anchor_one_pager(prs, idx=3, label="小鹏", brand="小鹏汽车",
        title="上海智能驾驶研发飞地 + 智驾出行运营基地",
        why_anchor="城市 NOA / XNGP 量产领先，急需上海「研发飞地 + 高端用户体验」。\n上海是新能源汽车销售第一大市，但小鹏华东研发尚未成型。",
        what_lacks="上海研发飞地（招聘 AI 人才）\n高端用户接待与品牌展示\nRobo 网约车运营试点\n数据闭环成本压缩",
        unique_offer=[
            ("B 栋 5,000 ㎡ 高端办公", "+ 静安人才公寓 80 套"),
            ("E 栋 1F 联合体验店", "中央广场城市发布"),
            ("60 个智驾车专用车位", "+ 户外测试区"),
            ("C 栋共享数据/仿真", "节省自建成本 30%"),
        ],
        deal_terms=[
            "6F + 7F 算法层共 ~3,366 ㎡ · 签约 5 年",
            "起始租金 7.5 元/㎡·天（算法软件价）",
            "免租 8 个月 · 装补 500 元/㎡",
            "9 月发布会主旨演讲位 + 媒体首发权",
            "1F 大堂咖啡角联合体验区（可冠名）",
        ],
        next_steps=[
            ("T+15d", "项目总监赴广州拜访李力耘"),
            ("T+45d", "邀请何小鹏出席 9 月发布会"),
            ("T+90d", "意向书签署"),
        ])

    anchor_one_pager(prs, idx=4, label="地平线", brand="地平线",
        title="静安「征程生态」开放实验室与算法-硬件协同中心",
        why_anchor="中国车规算力领跑，「征程」系列已成行业标配。\n急需上海生态拓展窗口 + 与算法/Tier1/车厂联调的实体场所。",
        what_lacks="算法-硬件联调场所（中心城区）\n城市 NOA 客户演示路段\nTier1 客户接待\n投资协同与上市资源",
        unique_offer=[
            ("C 栋 1,500 ㎡ 联合实验室", "「征程开放实验室」共建命名权"),
            ("户外测试 + V2X", "城市道缩比环线 + 路侧设备"),
            ("A 栋客户接待中心", "共享权 + 大客户演示"),
            ("产业基金协同", "区基金 + 冠松基金合计 5 亿优先 LP"),
        ],
        deal_terms=[
            "3F 核心研发整层 ~1,683 ㎡ · 净高 5.7 m",
            "起始租金 7.2 元/㎡·天",
            "免租 10 个月 · 装补 700 元/㎡",
            "「征程开放实验室」共建 + 优先调度",
            "上海科技新人才落户配额 20 个",
        ],
        next_steps=[
            ("T+10d", "邀请上海办参观户外测试区"),
            ("T+30d", "余凯（或黄畅）静安考察"),
            ("T+90d", "联合实验室共建协议"),
        ])

    anchor_one_pager(prs, idx=5, label="Momenta", brand="Momenta",
        title="上海「飞轮总部」 — 量产 + 数据闭环 + 出行三位一体基地",
        why_anchor="「飞轮」模式量产化领先，已与多家车企量产合作。\n上海有研发但分散，需要「总部级」叙事整合。",
        what_lacks="上海「总部级」门面\n量产数据-后市场闭环\n城市级 NOA 真实路测\n政府关系一企一策",
        unique_offer=[
            ("A 栋下半区或 B 栋整层", "冠名权 5 年"),
            ("冠松后市场", "4S/保险数据合规通道"),
            ("户外封闭 + 延伸路段", "城市级 NOA 真实路测"),
            ("静安区一企一策", "区领导班子直接对接"),
        ],
        deal_terms=[
            "8F 整层 ~1,683 ㎡ · 签约 6 年（含可叠加 9F）",
            "起始租金 5.8 元/㎡·天",
            "免租 14 个月 · 装补 1,000 元/㎡",
            "「Momenta 飞轮 · GS 数据沙盒」联合立项",
            "区级人才公寓 60 套 + 落户绿通",
        ],
        next_steps=[
            ("T+10d", "曹 CEO + 集团董事长 + 静安区长会面"),
            ("T+45d", "Term Sheet 草案"),
            ("T+120d", "签约"),
        ])

    # ============ 20. 300 家生态漏斗 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 2 · 任务 4",
               page_title="生态企业招商漏斗 · 300 家库",
               subtitle="5 家中介 · 9 类来源 · 端到端转化 17% · 60 家年签约目标")
    funnel = [
        ("L1 线索 (Lead)", 360, 12.0, BLUE),
        ("L2 触达 (Reached)", 290, 9.7, NAVY),
        ("L3 意向 (Intent)", 145, 4.8, BLUE),
        ("L4 谈判 (Negotiation)", 87, 2.9, NAVY),
        ("L5 签约 (Signed)", 60, 2.0, GOLD),
        ("L6 入驻 (Move-in)", 57, 1.9, GREEN),
    ]
    fy = Inches(1.25)
    for i, (lbl, n, w_inch, c) in enumerate(funnel):
        y = fy + Inches(0.65) * i
        add_rect(s, Inches(0.5), y, Inches(3.2), Inches(0.55), fill=CLOUD, line=LINE)
        add_text(s, Inches(0.55), y, Inches(3.1), Inches(0.55), lbl,
                 size=11, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(3.8), y, Inches(w_inch * 0.55), Inches(0.55), fill=c)
        add_text(s, Inches(3.8), y, Inches(w_inch * 0.55) + Inches(1.2), Inches(0.55),
                 f"  {n} 家", size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.5), Inches(5.5), Inches(10), Inches(0.4),
             "理论端到端转化（线索 → 签约）： 80% × 50% × 60% × 70% ≈ 17%",
             size=12, bold=True, color=NAVY)
    add_text(s, Inches(0.5), Inches(5.92), Inches(10), Inches(0.4),
             "60 家签约 ⇐ 360 有效线索 ⇐ 300 库 + 自拓溢出",
             size=12, color=BLUE)
    header = ["中介", "侧重", "佣金档", "年度任务"]
    rows = [
        ["戴德梁行 C&W", "跨国/外资 ADAS", "首月 100% / 长租 120%", "≥ 60 家"],
        ["仲量联行 JLL", "互联网/科技搬迁", "同上", "≥ 60 家"],
        ["高力 Colliers", "C+ 腰部企业", "同上", "≥ 50 家"],
        ["世邦 CBRE", "链主 / 整层", "链主最高 150%", "≥ 40 家"],
        ["本地精品行", "中小生态/工位", "首月 80%", "≥ 90 家"],
    ]
    add_text(s, Inches(8.2), Inches(1.0), Inches(5), Inches(0.35),
             "5 家中介渠道 · 非独家 + 30 天首报", size=13, bold=True, color=NAVY)
    add_table(s, Inches(8.2), Inches(1.4), Inches(4.8), Inches(3.1),
              header, rows,
              col_widths=[Inches(1.6), Inches(1.4), Inches(1.0), Inches(0.8)])
    add_text(s, Inches(8.2), Inches(4.7), Inches(5), Inches(0.35),
             "300 家库 · 产业链层级配比", size=13, bold=True, color=NAVY)
    bars = [("L4 算法", 60), ("L2 传感器", 50), ("L3 域控", 35), ("仿真/数据", 35),
            ("后市场", 30), ("投资/法律", 30), ("高校/科研", 30), ("L5 整车", 15),
            ("测试认证", 15)]
    by0 = Inches(5.05)
    for i, (lbl, n) in enumerate(bars):
        y = by0 + Inches(0.20) * i
        add_text(s, Inches(8.2), y, Inches(1.2), Inches(0.18), lbl, size=9, color=INK,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(9.5), y + Inches(0.04), Inches(n / 60.0 * 2.4), Inches(0.12), fill=BLUE)
        add_text(s, Inches(12.0), y, Inches(0.8), Inches(0.18), str(n), size=9, color=NAVY,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)

    # ============ 21. 政府关系 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 2 · 任务 5",
               page_title="政府关系对接 · 路径地图与政策包",
               subtitle="区/市双层 · 7 步首轮汇报 · 形成「一企一策」")
    add_text(s, Inches(0.5), Inches(1.15), Inches(6), Inches(0.4),
             "对接路径地图（市/区双层）", size=14, bold=True, color=NAVY)
    add_rect(s, Inches(0.5), Inches(1.55), Inches(6.2), Inches(5.3), fill=CLOUD, line=LINE)
    nodes = [
        ("市委市政府", Inches(2.5), Inches(1.7), Inches(2.0), Inches(0.45), NAVY, WHITE),
        ("市经信委", Inches(0.7), Inches(2.5), Inches(1.6), Inches(0.45), BLUE, WHITE),
        ("市公安/交通委", Inches(2.5), Inches(2.5), Inches(2.0), Inches(0.45), BLUE, WHITE),
        ("市科委/发改", Inches(4.7), Inches(2.5), Inches(1.8), Inches(0.45), BLUE, WHITE),
        ("静安区委区政府", Inches(2.3), Inches(3.5), Inches(2.4), Inches(0.5), GOLD, NAVY),
        ("区投促办", Inches(0.7), Inches(4.4), Inches(1.6), Inches(0.45), NAVY, WHITE),
        ("区经委/科委", Inches(2.5), Inches(4.4), Inches(2.0), Inches(0.45), NAVY, WHITE),
        ("区财政/税务", Inches(4.7), Inches(4.4), Inches(1.8), Inches(0.45), NAVY, WHITE),
        ("区交警支队", Inches(0.7), Inches(5.1), Inches(1.6), Inches(0.45), BLUE, WHITE),
        ("区人社/房管", Inches(2.5), Inches(5.1), Inches(2.0), Inches(0.45), BLUE, WHITE),
        ("区国资/街道", Inches(4.7), Inches(5.1), Inches(1.8), Inches(0.45), BLUE, WHITE),
    ]
    for txt, x, y, w, h, fc, tc in nodes:
        add_round(s, x, y, w, h, txt, fill=fc, color=tc, size=10, bold=True)
    add_text(s, Inches(0.7), Inches(5.9), Inches(6), Inches(0.4),
             "7 步首轮汇报 · 30 天完成全覆盖", size=12, bold=True, color=NAVY)
    add_text(s, Inches(0.7), Inches(6.25), Inches(6), Inches(0.6),
             "区投促办 → 副区长 → 区委书记/区长 → 市经信委 → 市公安交警 → 市科委/发改 → 区四套班子专题会",
             size=10, color=INK)
    add_text(s, Inches(7.0), Inches(1.15), Inches(6), Inches(0.4),
             "政策包 · 六维适配清单", size=14, bold=True, color=NAVY)
    pol = [
        ("牌照 / 测试", "封闭场地备案 · 1.5km 路测延伸 · L3/L4 试点联合体", BLUE),
        ("财税返还", "区级留成 80% 三年返 / 50% 后两年返", GOLD),
        ("人才 / 公寓", "链主 200 套 + 生态 300 套 · 落户绿通 80/年", NAVY),
        ("数据合规", "数交所沙盒 · 算力券最高 500 万 · 出境绿通", BLUE),
        ("产业基金", "区基金 1:1 配资 · 国调/中金/高瓴 LP 5–10 亿", GREEN),
        ("用地 / 配套", "测试区合规使用 · 周边路侧充电 · 一址多照", NAVY),
    ]
    for i, (t, d, c) in enumerate(pol):
        y = Inches(1.55) + Inches(0.85) * i
        add_rect(s, Inches(7.0), y, Inches(0.2), Inches(0.75), fill=c)
        add_rect(s, Inches(7.2), y, Inches(5.7), Inches(0.75), fill=CLOUD, line=LINE)
        add_text(s, Inches(7.35), y + Inches(0.05), Inches(5.3), Inches(0.32),
                 t, size=12, bold=True, color=NAVY)
        add_text(s, Inches(7.35), y + Inches(0.36), Inches(5.3), Inches(0.4),
                 d, size=10, color=INK)

    # ============ 22. Phase 3 扉页 ============
    section_cover(prs, "PHASE 3", "品牌与活动",
                  "Brand & Events · Make Noise, Build Trust",
                  ["任务 6 · 9 月旗舰发布会 + 年度 10 场活动 + 媒体清单"])

    # ============ 23. 9 月发布会 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 3 · 任务 6",
               page_title="9 月旗舰发布会 · 「中心城区 · 智驾未来」",
               subtitle="200 人现场 + 200 人在线 · 政企签约 · 政策包发布")
    agenda = [
        ("13:00–13:50", "注册 / 户外测试区分批参观"),
        ("14:00–14:25", "开场 + 静安区领导致辞"),
        ("14:25–14:35", "市经信委致辞"),
        ("14:35–14:55", "主旨发布：GS · iDrive Hub 产业定位"),
        ("14:55–15:30", "链主主题演讲 ×2"),
        ("15:30–16:00", "重大签约仪式（5–8 家）"),
        ("16:00–16:20", "茶歇 + 户外实车演示（直播）"),
        ("16:20–17:30", "圆桌：城市 NOA / 数据合规 / 中心城区基础设施"),
        ("17:30–17:45", "政策包发布：「iDrive · 静安 10 条」"),
        ("17:45–18:00", "园区生态启动 + 集体合影"),
        ("18:30–20:30", "招待晚宴 + 闭门 1v1"),
    ]
    add_text(s, Inches(0.5), Inches(1.2), Inches(6), Inches(0.4),
             "议程（9 月第三周 周四 14:00–20:30）", size=14, bold=True, color=NAVY)
    for i, (t, d) in enumerate(agenda):
        y = Inches(1.65) + Inches(0.40) * i
        add_rect(s, Inches(0.5), y, Inches(1.5), Inches(0.32), fill=NAVY)
        add_text(s, Inches(0.5), y, Inches(1.5), Inches(0.32), t,
                 size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(2.05), y, Inches(4.55), Inches(0.32), fill=CLOUD, line=LINE)
        add_text(s, Inches(2.15), y, Inches(4.45), Inches(0.32), d,
                 size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(7.0), Inches(1.2), Inches(6), Inches(0.4),
             "受邀 200 人 · 结构", size=14, bold=True, color=NAVY)
    invitees = [("政府领导", 25), ("链主 CEO", 30), ("生态代表", 60), ("投资机构", 25),
                ("高校研究院", 15), ("中介渠道", 10), ("媒体", 25), ("内部", 10)]
    for i, (lbl, n) in enumerate(invitees):
        y = Inches(1.65) + Inches(0.32) * i
        add_text(s, Inches(7.0), y, Inches(1.6), Inches(0.28), lbl,
                 size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(8.7), y + Inches(0.05), Inches(n / 60.0 * 3.0), Inches(0.18), fill=GOLD)
        add_text(s, Inches(11.8), y, Inches(0.8), Inches(0.28), str(n),
                 size=10, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
    add_text(s, Inches(7.0), Inches(4.3), Inches(6), Inches(0.4),
             "倒推时间表（T = 发布日）", size=14, bold=True, color=NAVY)
    timeline = ["T-90d 立项", "T-60d 物料一稿 / 邀请函", "T-45d 签约项目锁定",
                "T-30d 彩排 1", "T-21d 内容定稿", "T-14d 彩排 2 + 直播",
                "T-7d 彩排 3", "T-1d 全员彩排", "T 日 执行", "T+7d 复盘"]
    for i, t in enumerate(timeline):
        y = Inches(4.7) + Inches(0.22) * i
        add_round(s, Inches(7.0), y, Inches(0.3), Inches(0.18), "●",
                  fill=BLUE, color=WHITE, size=8, bold=True)
        add_text(s, Inches(7.4), y - Inches(0.01), Inches(5.5), Inches(0.22),
                 t, size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # ============ 24. 年度活动 + 媒体 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 3 · 任务 6",
               page_title="年度 10 场活动日历 + 媒体四层矩阵",
               subtitle="1 大会 + 4 季度沙龙 + 4 月度生态闭门会 + 1 国际峰会")
    header = ["#", "月份", "活动", "形式 / 规模", "主题方向"]
    rows = [
        ["1", "2026-05", "园区品牌发布会（启动）", "现场 80 人", "项目首秀 / 链主候选官宣"],
        ["2", "2026-06", "城市 NOA 商业化沙龙", "闭门 30 人", "链主 + 算法 + 出行"],
        ["3", "2026-07", "数据合规与算力训练营", "工作坊 50 人", "数据沙盒 + 算力券"],
        ["4", "2026-08", "智驾后市场闭门会（冠松特色）", "闭门 40 人", "保险 / 改装 / 售后"],
        ["5", "2026-09", "GS · iDrive Hub 全球招商发布会", "200 + 在线 200", "政企签约 / 政策包"],
        ["6", "2026-10", "仿真与端到端模型沙龙", "半日 60 人", "仿真生态 / 算力联训"],
        ["7", "2026-11", "智驾投融资 Demo Day", "半日 80 人", "早期算法 / 域控"],
        ["8", "2026-12", "年终招商盘点 + 入驻晚宴", "晚宴 120 人", "全年成果 / 续约"],
        ["9", "2027-01", "国际智驾产业峰会", "现场 250 人", "跨国车企 / 外资 ADAS"],
        ["10", "2027-03", "春季产业开放日 + 测试区开放周", "系列 5 天", "公众/媒体/高校"],
    ]
    add_table(s, Inches(0.5), Inches(1.2), Inches(8.0), Inches(5.3),
              header, rows,
              col_widths=[Inches(0.4), Inches(1.0), Inches(2.6), Inches(1.7), Inches(2.3)])
    add_text(s, Inches(8.8), Inches(1.2), Inches(4.5), Inches(0.4),
             "媒体四层矩阵", size=14, bold=True, color=NAVY)
    layers2 = [
        ("权威 · 政策", "央视 / 人民日报 / 新华社 / 解放日报 / 上观", NAVY),
        ("财经 · 商业", "财新 / 一财 / 21 世纪 / 36Kr / 钛媒体 / 界面", BLUE),
        ("行业 · 垂类", "汽车之家 / 懂车帝 / 高工智能 / 焉知 / 36 个公众号", GOLD),
        ("国际 · 英文", "Reuters / Bloomberg / FT / Nikkei / TechCrunch", GREEN),
    ]
    for i, (t, d, c) in enumerate(layers2):
        y = Inches(1.65) + Inches(1.20) * i
        add_rect(s, Inches(8.8), y, Inches(4.5), Inches(1.05), fill=CLOUD, line=LINE)
        add_rect(s, Inches(8.8), y, Inches(0.18), Inches(1.05), fill=c)
        add_text(s, Inches(9.05), y + Inches(0.08), Inches(4.2), Inches(0.4),
                 t, size=13, bold=True, color=NAVY)
        add_text(s, Inches(9.05), y + Inches(0.45), Inches(4.2), Inches(0.6),
                 d, size=10, color=INK)

    # ============ 25. Phase 4 扉页 ============
    section_cover(prs, "PHASE 4", "商业条款",
                  "Commercial · 4 Tiers · Service Pack · 3-Year Model",
                  ["任务 7 · 四档收费 + 服务包 + 报价单 + 三年财务"])

    # ============ 26. 四档收费 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 4 · 任务 7",
               page_title="四档收费方案矩阵",
               subtitle="链主月费+政策返还 / 生态月费 / 共享月费+服务佣金 / 对赌")
    cards4 = [
        ("甲档", "链主总部 (8–9F)", "月费 + 政策返还", NAVY,
         ["起始租金 5.5–6.5 元/㎡·天", "整层 1,683 ㎡ / 双层 ~3,366 ㎡",
          "免租 12–18 个月 · 装补 800–1,200/㎡",
          "区级税收 80% 三年返",
          "楼宇冠名 + 1F 大堂 + 嘉定/临港免费会员"]),
        ("乙档", "核心研发 (3–4F)", "月费 + 实验室共建", PURPLE,
         ["起始租金 7.0–8.0 元/㎡·天", "净层高 5.4–5.7 m （硬科技）",
          "免租 6–10 个月 · 装补 500–800/㎡",
          "3F 联合实验室命名权",
          "楼面荷载加固至 ≥ 7.5 kN/㎡"]),
        ("丙档", "算法软件 (6–7F)", "月费", BLUE,
         ["起始租金 7.5–8.5 元/㎡·天", "≥ 300 ㎡ 单元 (6–10 家)",
          "免租 4–8 个月 · 装补 300–500/㎡",
          "合同 3+3+1 年",
          "算力券 7 折 / 数据合规沙盒"]),
        ("丁档", "冠名 + 对赌", "低租 + 收益分成", GREEN,
         ["折扣租金 60%–70%", "1F+2F 冠名权 / 楼宇外立面",
          "营收/利润/估值对赌",
          "超额营收 5%–10% 返现",
          "园区基金 1%–3% 跟投权"]),
    ]
    cw = Inches(2.95); gx = Inches(0.15); cy = Inches(1.20); ch = Inches(5.4)
    for i, (lvl, name, sub, c, bullets) in enumerate(cards4):
        x = Inches(0.5) + (cw + gx) * i
        add_rect(s, x, cy, cw, ch, fill=CLOUD, line=LINE)
        add_rect(s, x, cy, cw, Inches(0.95), fill=c)
        add_text(s, x, cy + Inches(0.05), cw, Inches(0.4), lvl,
                 size=18, bold=True, color=GOLD if c == NAVY else WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, cy + Inches(0.40), cw, Inches(0.4), name,
                 size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x, cy + Inches(0.65), cw, Inches(0.3), sub,
                 size=10, color=WHITE, align=PP_ALIGN.CENTER)
        for j, b in enumerate(bullets):
            by = cy + Inches(1.10) + Inches(0.65) * j
            add_round(s, x + Inches(0.15), by + Inches(0.18), Inches(0.18), Inches(0.18), "●",
                      fill=c, color=WHITE, size=8, bold=True)
            add_text(s, x + Inches(0.40), by, cw - Inches(0.5), Inches(0.55), b,
                     size=10, color=INK)

    # ============ 27. 服务包矩阵 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 4 · 任务 7",
               page_title="服务包矩阵 · 把租金以外的价值显性化",
               subtitle="6 大类 · 12 项服务 · 按四档差异化配置 · 是租金等效折扣的核心抓手")
    sv_header = ["服务", "甲 链主总部", "乙 核心研发", "丙 算法软件", "丁 冠名+对赌"]
    sv_rows = [
        ["1F 大堂冠名", "默认 5 年", "—", "—", "默认（叠加 2F 展厅）"],
        ["楼宇外立面冠名", "可选", "—", "—", "默认"],
        ["屋顶花园专属", "默认", "—", "—", "默认"],
        ["联合实验室共建", "—", "命名权", "—", "命名权"],
        ["嘉定/临港测试场会员代办", "免费", "8 折", "9 折", "免费"],
        ["路测延伸路段协调", "优先", "平等", "平等", "优先"],
        ["共享算力券", "50 万", "30 万", "7 折", "80 万"],
        ["数据合规沙盒", "白名单 + 优先", "标准", "联合申请", "白名单 + 优先"],
        ["招聘联运 / 落户绿通", "50 个/年", "15 个/年", "集体面试", "30 个/年"],
        ["政策代办（一窗通办）", "一企一策专人", "标准代办", "自助 + 模板", "一企一策专人"],
        ["公关与媒体", "主旨发布 + 首发", "主题分享", "Demo 抽签", "主旨发布 + 首发"],
        ["后市场协同 (冠松)", "数据接口", "应用试点", "—", "数据接口"],
    ]
    add_table(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5),
              sv_header, sv_rows,
              col_widths=[Inches(2.6), Inches(2.4), Inches(2.4), Inches(2.4), Inches(2.5)],
              header_size=11, body_size=9, body_align=PP_ALIGN.LEFT)
    add_text(s, Inches(0.5), Inches(6.85), Inches(12), Inches(0.3),
             "把租金以外的服务量化为『积分券』——为客户算清账，为我们沉淀粘性",
             size=10, italic=True, color=GREY)

    # ============ 28. 链主报价单样张 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 4 · 任务 7",
               page_title="链主总部报价单样张（甲档 · 保密）",
               subtitle="致 [链主]：8F + 9F 整层 ~3,366 ㎡ · 6 年 · 5.8 元/㎡·天 · 一企一策政策包")
    # 标题与抬头
    add_rect(s, Inches(0.5), Inches(1.20), Inches(12.3), Inches(0.55), fill=NAVY)
    add_text(s, Inches(0.7), Inches(1.20), Inches(8), Inches(0.55),
             "GS · iDrive Hub  |  链主独栋报价单  |  CONFIDENTIAL",
             size=14, bold=True, color=GOLD, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(8.5), Inches(1.20), Inches(4.0), Inches(0.55),
             "有效期：30 日", size=11, color=WHITE,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    # 左列：商业条款
    add_text(s, Inches(0.5), Inches(1.95), Inches(6), Inches(0.4),
             "一、商业条款", size=13, bold=True, color=NAVY)
    biz = [
        ("标的", "01# 8F + 9F 整层 ~3,366 ㎡（4.2 / 4.3 m 净高）"),
        ("租赁期", "6 年（含免租期）"),
        ("起始租金", "5.8 元/㎡·天 · 物业 28 元/㎡·月"),
        ("免租期", "15 个月"),
        ("调价机制", "3 年一调 · 调幅 = max(CPI, 5%)"),
        ("装补", "1,000 元/㎡ · 封顶 336 万元 · 里程碑分期"),
        ("履约保证", "6 个月租金"),
        ("续约", "到期前 12 个月可优先续约 5 年"),
    ]
    for i, (k, v) in enumerate(biz):
        y = Inches(2.40) + Inches(0.36) * i
        add_round(s, Inches(0.55), y + Inches(0.06), Inches(1.3), Inches(0.24),
                  k, fill=BLUE, color=WHITE, size=10, bold=True)
        add_text(s, Inches(1.95), y, Inches(4.7), Inches(0.36), v,
                 size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # 中列：政策包
    add_text(s, Inches(6.85), Inches(1.95), Inches(6), Inches(0.4),
             "二、一企一策政策包", size=13, bold=True, color=NAVY)
    pol2 = [
        ("税收返还", "区级留成 80% 三年返 / 50% 后两年返"),
        ("一次性奖励", "[一事一议]"),
        ("人才公寓", "80–120 套"),
        ("落户绿通", "高级研发 20 个 / 年"),
        ("测试场会员", "嘉定/临港封闭场代办（终身免费）"),
        ("路测延伸", "1.5 km 静安区路测专项联合申请"),
        ("数据合规", "上海数交所沙盒白名单"),
        ("品牌权益", "1F 大堂 + 楼宇外立面冠名 5 年 + 9 月主旨发布"),
    ]
    for i, (k, v) in enumerate(pol2):
        y = Inches(2.40) + Inches(0.36) * i
        add_round(s, Inches(6.90), y + Inches(0.06), Inches(1.3), Inches(0.24),
                  k, fill=GOLD, color=NAVY, size=10, bold=True)
        add_text(s, Inches(8.30), y, Inches(4.4), Inches(0.36), v,
                 size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # 底部：签约里程碑 + 联系
    add_rect(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.5), fill=BLUE)
    add_text(s, Inches(0.7), Inches(5.4), Inches(12), Inches(0.5),
             "三、签约里程碑：T+30d 意向条款书 → T+90d 正式合同 → T+150d 进场装修 → T+330d 入驻 + 区领导剪彩",
             size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.95), fill=NAVY)
    add_text(s, Inches(0.7), Inches(6.05), Inches(12), Inches(0.4),
             "冠松集团 · GS · iDrive Hub 项目组",
             size=12, bold=True, color=GOLD)
    add_text(s, Inches(0.7), Inches(6.45), Inches(12), Inches(0.45),
             "联系人：项目总监 [姓名]  |  电话：[xxx]  |  邮箱：[xxx@xxx.com]",
             size=11, color=WHITE)

    # ============ 29. 三年财务测算 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 4 · 任务 7",
               page_title="三年财务测算 · 关键 KPI",
               subtitle="基于 1.5 万方 实测口径 · 出租净 8,300 ㎡ · Y3 EBITDA 转正约 1,400 万元")
    header = ["科目（万元）", "Y1 2026", "Y2 2027", "Y3 2028"]
    rows = [
        ["租金收入", "615", "1,442", "2,091"],
        ["物业费收入", "97", "195", "257"],
        ["冠名 / 服务 / 后市场", "230", "700", "1,400"],
        ["政策返还（净计入）", "80", "600", "1,500"],
        ["基金管理费", "0", "300", "500"],
        ["总收入", "1,022", "3,237", "5,748"],
        ["总成本", "3,130", "3,790", "4,380"],
        ["EBITDA", "−2,108", "−553", "+1,368"],
        ["EBIT", "−3,308", "−1,753", "+168"],
        ["税前利润", "−3,908", "−2,333", "−352"],
    ]
    add_table(s, Inches(0.5), Inches(1.2), Inches(7.0), Inches(4.5),
              header, rows,
              col_widths=[Inches(2.4), Inches(1.5), Inches(1.5), Inches(1.6)])
    kpis = [
        ("入驻率", ["35%", "70%", "92%"], BLUE),
        ("链主签约（累计）", ["0", "1", "1"], NAVY),
        ("入驻企业（累计）", ["4", "9", "14"], GOLD),
        ("EBITDA 利润率", ["−206%", "−17%", "+24%"], GREEN),
    ]
    add_text(s, Inches(7.8), Inches(1.2), Inches(5.5), Inches(0.4),
             "三年关键 KPI", size=14, bold=True, color=NAVY)
    for i, (lbl, vals, c) in enumerate(kpis):
        y = Inches(1.65) + Inches(0.95) * i
        add_text(s, Inches(7.8), y + Inches(0.05), Inches(2.3), Inches(0.4),
                 lbl, size=12, bold=True, color=NAVY)
        for j, v in enumerate(vals):
            x = Inches(10.1) + Inches(1.0) * j
            add_round(s, x, y + Inches(0.02), Inches(0.85), Inches(0.55),
                      v, fill=c, color=WHITE, size=12, bold=True)
            add_text(s, x, y + Inches(0.6), Inches(0.85), Inches(0.25),
                     f"Y{j+1}", size=8, color=GREY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.85), Inches(13), Inches(0.4),
             "敏感性（Y3 EBITDA 万元变化）：入驻 ±10% → −480/+220 · 租金 ±10% → ±210 · 政策兑现 −30% → −450 · 冠名/服务 −30% → −420",
             size=11, color=INK)

    # ============ NEW 29b. 服务平台佣金构成 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 4 · 任务 7 补充",
               page_title="服务平台佣金 600 万 = 议价权变现",
               subtitle="把入驻企业聚合后的议价权变现 · 6 个子项构成 · 增量人力成本 ≤ 100 万")

    # 左：6 个子项卡片（金字塔结构）
    add_text(s, Inches(0.5), Inches(1.15), Inches(6), Inches(0.4),
             "Y3 600 万构成（6 大子项）", size=14, bold=True, color=NAVY)
    items_svc = [
        ("A. 算力服务转售", "250 万",
         "阿里云 / 华为云 / 上海超算 · 价差 8–15%", BLUE),
        ("B. 政府事务服务费", "150 万",
         "一企一策 / 牌照代办 / 落户 · 单项 5–20 万", GOLD),
        ("C. 联合实验室分成", "80 万",
         "3F 实验室对外承接 Tier1 · 出资比例分成", PURPLE),
        ("D. 招聘联运", "60 万",
         "高级研发岗位渠道分成 · 单笔 8–15%", GREEN),
        ("E. 法务/IP/咨询", "60 万",
         "君合/方达/毕马威 驻点 · 单笔 5–10%", NAVY),
    ]
    for i, (t, amount, d, c) in enumerate(items_svc):
        y = Inches(1.55) + Inches(0.92) * i
        add_rect(s, Inches(0.5), y, Inches(0.18), Inches(0.82), fill=c)
        add_rect(s, Inches(0.7), y, Inches(6.0), Inches(0.82), fill=CLOUD, line=LINE)
        add_text(s, Inches(0.85), y + Inches(0.05), Inches(2.2), Inches(0.35),
                 t, size=12, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_round(s, Inches(3.10), y + Inches(0.20), Inches(0.95), Inches(0.40),
                  amount, fill=c, color=WHITE if c != GOLD else NAVY,
                  size=12, bold=True)
        add_text(s, Inches(4.20), y + Inches(0.05), Inches(2.4), Inches(0.73),
                 d, size=9, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # 右上：三年爬坡柱状
    add_text(s, Inches(7.0), Inches(1.15), Inches(6), Inches(0.4),
             "三年爬坡（与入驻率挂钩 70%）", size=14, bold=True, color=NAVY)
    bar_max = 600
    bar_w_max = Inches(5.5)
    bar_data = [("Y1 (35%)", 80, BLUE), ("Y2 (70%)", 250, GOLD), ("Y3 (92%)", 600, GREEN)]
    for i, (lbl, n, c) in enumerate(bar_data):
        y = Inches(1.55) + Inches(0.7) * i
        add_text(s, Inches(7.0), y, Inches(1.4), Inches(0.50),
                 lbl, size=12, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        bw = Inches(n / bar_max * 5.5)
        add_rect(s, Inches(8.4), y + Inches(0.10), bw, Inches(0.30), fill=c)
        add_text(s, Inches(8.4) + bw + Inches(0.05), y + Inches(0.05),
                 Inches(1.5), Inches(0.40),
                 f"{n} 万", size=12, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)

    # 右下：核心特点
    add_rect(s, Inches(7.0), Inches(4.0), Inches(6.0), Inches(2.3), fill=NAVY)
    add_text(s, Inches(7.15), Inches(4.05), Inches(5.7), Inches(0.4),
             "★ 商业模式特点", size=13, bold=True, color=GOLD)
    bullets_svc = [
        "议价权变现：1 链主 + 14 家生态体量议价",
        "边际成本低：增量人力成本 ≤ 100 万/年",
        "强生态依附：70% 与入驻量挂钩，弹性大",
        "政策合规：不承诺审批结果 / 不收回扣",
    ]
    for i, b in enumerate(bullets_svc):
        y = Inches(4.50) + Inches(0.42) * i
        add_round(s, Inches(7.15), y + Inches(0.10), Inches(0.20), Inches(0.20),
                  "●", fill=GOLD, color=NAVY, size=8, bold=True)
        add_text(s, Inches(7.45), y, Inches(5.4), Inches(0.4),
                 b, size=11, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    # 底部里程碑
    add_text(s, Inches(0.5), Inches(6.45), Inches(12.5), Inches(0.40),
             "实操里程碑：M3 算力供应商签约 · M5 法务/会计师驻点 · M8 联合实验室开放 · M12 月度 GMV ≥ 100 万",
             size=11, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)

    # ============ NEW 29c. 后市场协同分成构成 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 4 · 任务 7 补充",
               page_title="后市场协同分成 300 万 = 冠松独家壁垒变现",
               subtitle="冠松 6 大资源 × 园区智驾产品 · 4 道不可复制护城河")

    # 左：6 个子项
    add_text(s, Inches(0.5), Inches(1.15), Inches(6), Inches(0.4),
             "Y3 300 万构成（6 大子项）", size=14, bold=True, color=NAVY)
    items_amk = [
        ("A. 数据闭环分成", "120 万",
         "4S 真实数据→数交所沙盒→链主订阅 · 抽 10–15%", NAVY),
        ("B. 保险定损协同", "80 万",
         "智驾保险产品 × 冠松定损 · 与平安/人保/太保分成", BLUE),
        ("C. 二手智驾车认证", "30 万",
         "残值评估 + 智驾包激活流转", GOLD),
        ("D. 测试车队融资", "30 万",
         "冠松融资租赁子公司利率低 1.5–2%", PURPLE),
        ("E. 冠松车队数据采集", "30 万",
         "100+ 营运车 · 真实上海典型场景", GREEN),
        ("F. 体验店流量分成", "10 万",
         "1F 大堂体验店 · 园区企业产品分销", BLUE),
    ]
    for i, (t, amount, d, c) in enumerate(items_amk):
        y = Inches(1.55) + Inches(0.78) * i
        add_rect(s, Inches(0.5), y, Inches(0.18), Inches(0.70), fill=c)
        add_rect(s, Inches(0.7), y, Inches(6.0), Inches(0.70), fill=CLOUD, line=LINE)
        add_text(s, Inches(0.85), y + Inches(0.03), Inches(2.2), Inches(0.30),
                 t, size=11, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_round(s, Inches(3.10), y + Inches(0.15), Inches(0.95), Inches(0.40),
                  amount, fill=c, color=WHITE if c != GOLD else NAVY,
                  size=11, bold=True)
        add_text(s, Inches(4.20), y + Inches(0.03), Inches(2.4), Inches(0.63),
                 d, size=9, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # 右上：三年爬坡
    add_text(s, Inches(7.0), Inches(1.15), Inches(6), Inches(0.4),
             "三年爬坡（依赖链主入驻）", size=14, bold=True, color=NAVY)
    bar_data2 = [("Y1", 50, BLUE), ("Y2 (链主 1)", 150, GOLD), ("Y3 (链主 1+生态 13)", 300, GREEN)]
    for i, (lbl, n, c) in enumerate(bar_data2):
        y = Inches(1.55) + Inches(0.55) * i
        add_text(s, Inches(7.0), y, Inches(1.8), Inches(0.40),
                 lbl, size=11, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        bw = Inches(n / 300.0 * 3.0)
        add_rect(s, Inches(8.8), y + Inches(0.08), bw, Inches(0.25), fill=c)
        add_text(s, Inches(8.8) + bw + Inches(0.05), y + Inches(0.05),
                 Inches(1.2), Inches(0.30),
                 f"{n} 万", size=11, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)

    # 右下：4 道护城河
    add_text(s, Inches(7.0), Inches(3.4), Inches(6), Inches(0.4),
             "4 道护城河（不可复制）", size=14, bold=True, color=NAVY)
    moats = [
        ("① 数据真实性", "4S 真实数据自然产生", NAVY),
        ("② 保险牌照", "持牌保险经纪/公估机构合作", BLUE),
        ("③ 资金成本", "融资租赁利率低市场 1.5–2%", GOLD),
        ("④ 数据合规", "数交所沙盒处理 · 链主可合规使用", GREEN),
    ]
    for i, (t, d, c) in enumerate(moats):
        y = Inches(3.85) + Inches(0.62) * i
        add_rect(s, Inches(7.0), y, Inches(0.18), Inches(0.55), fill=c)
        add_rect(s, Inches(7.2), y, Inches(5.8), Inches(0.55), fill=CLOUD, line=LINE)
        add_text(s, Inches(7.35), y + Inches(0.05), Inches(1.8), Inches(0.45),
                 t, size=11, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(9.15), y + Inches(0.05), Inches(3.8), Inches(0.45),
                 d, size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # 底部里程碑
    add_rect(s, Inches(0.5), Inches(6.40), Inches(12.3), Inches(0.45), fill=NAVY)
    add_text(s, Inches(0.7), Inches(6.40), Inches(12.0), Inches(0.45),
             "实操里程碑：M2 子公司框架协议 · M4 数交所沙盒 · M6 首笔数据订阅 · M9 智驾保险合作 · M12 测试车队融资",
             size=10, bold=True, color=GOLD, anchor=MSO_ANCHOR.MIDDLE)

    # ============ 30. Phase 5 扉页 ============
    section_cover(prs, "PHASE 5", "落地推进",
                  "Rollout · 12 Months · 4 → 22 Team",
                  ["任务 8 · 12 个月里程碑 + 团队扩编 + RACI + 风险矩阵"])

    # ============ 31. 12 个月甘特 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 5 · 任务 8",
               page_title="12 个月里程碑甘特图",
               subtitle="M1 团队就位 → M5 9 月发布会 → M12 入驻率 65%")
    months = [f"M{i+1}" for i in range(12)]
    chart_x = Inches(2.8); chart_y = Inches(1.2); chart_w = Inches(10.0)
    col_w = chart_w / 12.0
    for i, m in enumerate(months):
        x = chart_x + col_w * i
        add_rect(s, x, chart_y, col_w, Inches(0.35), fill=NAVY)
        add_text(s, x, chart_y, col_w, Inches(0.35), m,
                 size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tasks = [
        ("立项与法人主体",       0, 1, BLUE),
        ("团队组建（4→22）",     0, 12, NAVY),
        ("样板间 / 视觉",        1, 3, GOLD),
        ("户外测试区一期",       1, 5, GREEN),
        ("品牌官网 + CRM",       1, 3, BLUE),
        ("300 家库 + 中介签约",  2, 3, NAVY),
        ("首轮政府汇报",         1, 3, GOLD),
        ("政策包 v0 → v1",       3, 5, BLUE),
        ("链主 TOP5 接触",       1, 9, NAVY),
        ("链主首份 Term Sheet",  2, 2, GOLD),
        ("生态意向 30 家",       4, 4, BLUE),
        ("5 月品牌发布会",       4, 1, GREEN),
        ("9 月旗舰发布会",       8, 1, RED),
        ("「iDrive 静安 10 条」",8, 2, GOLD),
        ("入驻率 50%",           10, 1, GREEN),
        ("入驻率 65%",           11, 1, GREEN),
    ]
    row_h = Inches(0.28)
    for i, (lbl, start, dur, c) in enumerate(tasks):
        y = chart_y + Inches(0.4) + (row_h + Inches(0.03)) * i
        add_text(s, Inches(0.4), y, Inches(2.35), row_h, lbl,
                 size=9, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, chart_x, y, chart_w, row_h, fill=CLOUD)
        bx = chart_x + col_w * start
        bw = col_w * dur
        add_rect(s, bx + Emu(20000), y + Emu(20000), bw - Emu(40000), row_h - Emu(40000), fill=c)

    # ============ 32. 团队 + RACI ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 5 · 任务 8",
               page_title="4 人核心团队 → 22 人稳态 + RACI 分工",
               subtitle="启动 M1–3：4 人 / 扩张 M4–6：14 人 / 稳态 M7–12：22 人")
    add_text(s, Inches(0.5), Inches(1.15), Inches(6), Inches(0.4),
             "4 人核心团队（启动期）", size=14, bold=True, color=NAVY)
    roles = [
        ("项目总监 GM", "全面统筹 / 链主谈判 / 重大客户", NAVY),
        ("GR 总监", "政府关系 / 政策包 / 合规", BLUE),
        ("招商总监", "渠道 / 漏斗 / 生态签约", GOLD),
        ("运营总监", "物业 / IT / 测试区 / 活动", GREEN),
    ]
    for i, (t, d, c) in enumerate(roles):
        y = Inches(1.55) + Inches(0.95) * i
        add_rect(s, Inches(0.5), y, Inches(6.0), Inches(0.85), fill=CLOUD, line=LINE)
        add_rect(s, Inches(0.5), y, Inches(0.18), Inches(0.85), fill=c)
        add_text(s, Inches(0.7), y + Inches(0.08), Inches(5.7), Inches(0.4),
                 t, size=14, bold=True, color=NAVY)
        add_text(s, Inches(0.7), y + Inches(0.45), Inches(5.7), Inches(0.4),
                 d, size=11, color=INK)
    add_text(s, Inches(0.5), Inches(5.6), Inches(6), Inches(0.4),
             "扩编节奏：4 → 14 → 22 人", size=12, bold=True, color=NAVY)
    add_text(s, Inches(0.5), Inches(5.95), Inches(6), Inches(0.6),
             "招商 1→5→8 · GR 1→2→2 · 运营 1→3→5 · 市场 0→1→3 · 法务/财务/HR 0→3→3",
             size=10, color=INK)
    add_text(s, Inches(7.0), Inches(1.15), Inches(6), Inches(0.4),
             "RACI 关键事项分工", size=14, bold=True, color=NAVY)
    header = ["关键事项", "项目总监", "GR", "招商", "运营", "董事长", "法务"]
    rows = [
        ["链主谈判", "R", "C", "C", "I", "A", "C"],
        ["政府汇报 / 政策包", "C", "R", "I", "I", "A", "R"],
        ["中介渠道 / 漏斗例会", "A", "I", "R", "I", "I", "C"],
        ["9 月发布会 / 测试区", "A", "C", "C", "R", "C", "I"],
        ["合同 / 对赌", "A", "C", "C", "I", "A", "R"],
        ["财务模型 / 预算", "A", "I", "C", "C", "A", "I"],
        ["数据合规", "C", "C", "C", "C", "A", "R"],
    ]
    add_table(s, Inches(7.0), Inches(1.55), Inches(6.0), Inches(4.0),
              header, rows,
              col_widths=[Inches(2.0), Inches(0.7), Inches(0.6), Inches(0.7),
                          Inches(0.7), Inches(0.8), Inches(0.6)],
              header_size=10, body_size=10, body_align=PP_ALIGN.CENTER)
    add_text(s, Inches(7.0), Inches(5.65), Inches(6.0), Inches(0.5),
             "R = Responsible · A = Accountable · C = Consulted · I = Informed",
             size=10, italic=True, color=GREY)

    # ============ 33. 风险矩阵 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 5 · 任务 8",
               page_title="风险矩阵 · 高频高影响项的对冲方案",
               subtitle="按概率 × 影响排序，每项配对应『对冲手段』和『启动条件』")
    rk_header = ["#", "风险", "概率", "影响", "对冲手段", "启动条件"]
    rk_rows = [
        ["R1", "链主推进慢于预期", "中", "高", "启动备选名单（蔚来/理想/大众问问/Smart）+ 链主分阶段意向", "TOP5 中 ≤1 进入 Term Sheet @ M6"],
        ["R2", "户外测试区报批延迟", "中", "高", "第三方代办 + 先以 V2X 试验段名义启用 + 与中汽研合作过渡", "M4 仍未取得备案"],
        ["R3", "政策返还兑现延迟", "中", "中", "改约定 + 财政确认函分期，避免谈判停滞 + 集团财政担保过渡", "区财政书面确认延误 > 60 天"],
        ["R4", "媒体或事故敏感期", "低", "高", "公关年框 SOP，重大突发暂停所有路测对外发声 + 法务声明库", "行业重大事故 / 监管整顿"],
        ["R5", "团队建设速度", "中", "中", "启动期接受外包（PR/HR/法务）+ 快速 Convert · 头部猎头授权", "招聘缺口 > 30% @ M3"],
        ["R6", "资金压力", "中", "中", "Y2 起 REITs 化探索 + 集团信用借款 + 工程款分期", "现金安全垫 < 6 个月"],
        ["R7", "竞品价格战", "中", "中", "服务包等效折扣 + 一企一策护城河 + 对赌档差异化", "嘉定/张江同期推出新政"],
        ["R8", "数据合规 / 出口管制", "低", "高", "上海数交所沙盒 + 法务审查 + 客户分级（外资单独通道）", "新法规出台 / 客户被审查"],
        ["R9", "施工 / 装修延期", "中", "中", "工程总包 + 关键节点违约金 + 分阶段交付", "里程碑滞后 > 30 天"],
        ["R10", "物业运营 / 安全事故", "低", "高", "Tier III 等同数据机房 + 24h 安保 + 保险年框", "首次重大事故立刻启动"],
    ]
    add_table(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5),
              rk_header, rk_rows,
              col_widths=[Inches(0.5), Inches(2.4), Inches(0.7), Inches(0.7),
                          Inches(4.5), Inches(3.5)],
              header_size=11, body_size=9)
    add_text(s, Inches(0.5), Inches(6.85), Inches(13), Inches(0.3),
             "风险委员会按月评审；任何风险升至「红」状态，72 小时内升级至集团董事长 + 项目总监 + 法务三方专项会",
             size=10, italic=True, color=GREY)

    # ====================================================================
    # ============ v1.2 新增：谈判 / 团队 / 资源 / KPI / FAQ / 索引 =======
    # ====================================================================

    # ============ 新 1. 谈判策略总览 · 四类对手 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="附 · 谈判策略",
               page_title="谈判策略总览 · 四类对手画像",
               subtitle="链主 / 政府 / 中介 / 生态客户 · 不同节奏 · 不同主谈 · 统一原则")
    quads_n = [
        ("I · 链主", "CEO / 总裁 / CTO + 政府事务 + 法务\n\n主谈：项目总监 + 集团董事长\n备谈：GR + 法务\n\n节奏：30–90 天 5 轮\n（建联 → 踏勘 → TS → 合同 → 落定）",
         NAVY),
        ("II · 政府", "区委书记 / 区长 / 副区长 / 区投促办 / 市经信委\n\n主谈：GR 总监 + 集团董事长\n备谈：项目总监\n\n节奏：7–30 天 3 轮\n（投促办 → 副区长 → 四套班子）",
         BLUE),
        ("III · 中介", "经纪团队 GM / Sr. Director\n\n主谈：招商总监\n备谈：销售经理 + 法务\n\n节奏：7–14 天 2 轮\n（议价 → 签约 + 启动会）",
         GOLD),
        ("IV · 生态客户", "创始人 / CEO / HR / 行政\n\n主谈：招商经理\n备谈：销售经理 + 客户成功\n\n节奏：14–30 天 3 轮\n（电话 → 踏勘 → 合同）",
         GREEN),
    ]
    qw = Inches(6.15); qh = Inches(2.85)
    positions = [(Inches(0.5), Inches(1.20)), (Inches(6.85), Inches(1.20)),
                 (Inches(0.5), Inches(4.20)), (Inches(6.85), Inches(4.20))]
    for (px, py), (t, d, c) in zip(positions, quads_n):
        add_rect(s, px, py, qw, qh, fill=CLOUD, line=LINE)
        add_rect(s, px, py, qw, Inches(0.55), fill=c)
        add_text(s, px + Inches(0.15), py + Inches(0.05), qw - Inches(0.3), Inches(0.45),
                 t, size=18, bold=True, color=GOLD if c == NAVY else WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, px + Inches(0.20), py + Inches(0.65), qw - Inches(0.4), qh - Inches(0.75),
                 d, size=11, color=INK)

    # ============ 新 2. 链主谈判 · 5 轮节奏 + 三层底线 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="附 · 谈判策略",
               page_title="链主谈判 · 5 轮节奏 + 三层底线",
               subtitle="谈判前必须内部明确底线（绿/黄/红）· 超线立即休会请示")
    # 左：5 轮节奏
    add_text(s, Inches(0.5), Inches(1.15), Inches(6), Inches(0.4),
             "5 轮谈判节奏（30–90 天）", size=14, bold=True, color=NAVY)
    rounds = [
        ("R1 建联", "T+0", "项目总监送一页纸提案", BLUE),
        ("R2 踏勘", "T+15d", "现场踏勘 + 报价单", BLUE),
        ("R3 议定", "T+45d", "Term Sheet + 财务对比", GOLD),
        ("R4 合同", "T+75d", "合同条款逐条 + 法务", PURPLE),
        ("R5 落定", "T+90d", "签约 + 区领导见证", GREEN),
    ]
    for i, (r, t, d, c) in enumerate(rounds):
        y = Inches(1.55) + Inches(1.0) * i
        add_rect(s, Inches(0.5), y, Inches(0.85), Inches(0.85), fill=c)
        add_text(s, Inches(0.5), y, Inches(0.85), Inches(0.4), r,
                 size=11, bold=True, color=NAVY if c == GOLD else WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(0.5), y + Inches(0.42), Inches(0.85), Inches(0.4), t,
                 size=10, bold=True, color=NAVY if c == GOLD else WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(1.40), y, Inches(5.1), Inches(0.85), fill=CLOUD, line=LINE)
        add_text(s, Inches(1.55), y, Inches(4.95), Inches(0.85),
                 d, size=12, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # 右：三层底线
    add_text(s, Inches(7.0), Inches(1.15), Inches(6), Inches(0.4),
             "三层底线（红/黄/绿）", size=14, bold=True, color=NAVY)
    bottom_header = ["条款", "🟢 理想", "🟡 目标", "🔴 底线"]
    bottom_rows = [
        ["起始租金 元/㎡·天", "6.5", "5.8", "5.0"],
        ["免租期 (月)", "9", "15", "24"],
        ["装补 元/㎡", "600", "1,000", "1,500"],
        ["合同期 (年)", "8", "6", "5"],
        ["履约保证 (月)", "6", "6", "3"],
        ["调价机制", "每3y max(CPI,5%)", "每3y max(CPI,4%)", "每5y max(CPI,3%)"],
        ["政策返还", "80% 三年", "80% 三年+50%两年", "80% 三年+30%两年"],
        ["楼宇冠名 (年)", "5", "5", "3"],
    ]
    add_table(s, Inches(7.0), Inches(1.55), Inches(6.0), Inches(4.8),
              bottom_header, bottom_rows,
              col_widths=[Inches(2.0), Inches(1.3), Inches(1.4), Inches(1.3)],
              header_size=10, body_size=9, body_align=PP_ALIGN.CENTER)
    add_rect(s, Inches(7.0), Inches(6.50), Inches(6.0), Inches(0.40), fill=NAVY)
    add_text(s, Inches(7.15), Inches(6.50), Inches(5.7), Inches(0.40),
             "超过授权区间立即休会请示集团董事长 · 不在书面承诺政府未确认事项",
             size=10, bold=True, color=GOLD, anchor=MSO_ANCHOR.MIDDLE)

    # ============ 新 3. 链主让步阶梯（Concession Ladder） ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="附 · 谈判策略",
               page_title="链主让步阶梯 · 每让一步必有「换」",
               subtitle="5 步让步 · 对应 NPV 影响 · 严禁无条件让步")
    add_text(s, Inches(0.5), Inches(1.15), Inches(12), Inches(0.4),
             "让步阶梯（Concession Ladder）", size=14, bold=True, color=NAVY)
    add_rect(s, Inches(0.5), Inches(1.55), Inches(12.3), Emu(20000), fill=GOLD)

    ladder_header = ["步", "我方让步", "幅度", "NPV 影响 (万元)", "对方对等承诺", "授权层"]
    ladder_rows = [
        ["1", "起始租金 6.5 → 6.0 元/㎡·天", "−7.7%", "−550", "8 年长租 + 6 个月履约", "项目总监"],
        ["2", "免租期 12 → 15 个月", "+25%", "−240", "装补封顶 1,000/㎡（不浮动）", "项目总监"],
        ["3", "装补 800 → 1,000 元/㎡", "+25%", "−67", "楼宇冠名 5 年 + 9 月主旨", "项目总监"],
        ["4", "政策返还 70% → 80%", "+14%", "−340", "以政府专班书面文件兑现", "GR + 集团董事长"],
        ["5", "续约权（涨幅封顶 5%）", "—", "−90", "退租赔偿 ≥ 60% 政策返还", "项目总监"],
        ["🔴", "底线：5.0 / 24 / 1,500 / 5 年", "—", "−1,500+", "超线请示集团董事长", "集团董事长"],
    ]
    add_table(s, Inches(0.5), Inches(1.75), Inches(12.3), Inches(4.5),
              ladder_header, ladder_rows,
              col_widths=[Inches(0.6), Inches(3.8), Inches(1.0), Inches(1.8),
                          Inches(3.4), Inches(1.7)],
              header_size=11, body_size=10)

    add_rect(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.55), fill=NAVY)
    add_text(s, Inches(0.7), Inches(6.42), Inches(12.0), Inches(0.50),
             "★ 让步原则：每让一步必带条件 · 小步多次 · 让 5 次小步胜过让 1 次大步 · 24h 内出书面纪要",
             size=11, bold=True, color=GOLD, anchor=MSO_ANCHOR.MIDDLE)

    # ============ 新 4. 政府 / 中介 / 生态 谈判要点 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="附 · 谈判策略",
               page_title="政府 / 中介 / 生态 三类谈判要点",
               subtitle="不同对手不同打法 · 但都共享统一谈判工具箱")
    cards = [
        ("II · 政府谈判", "7–30 天 3 轮", NAVY,
         ["从「汇报」到「专班」",
          "从「政策包」到「一企一策」",
          "从「承诺」到「书面」",
          "从「返还」到「确认函」",
          "从「口袋书」到「政绩素材」"]),
        ("III · 中介谈判", "7–14 天 2 轮", BLUE,
         ["佣金 80% → 100% ⇄ 任务量 60 家",
          "链主 → 120% ⇄ 30 天首报独家",
          "续约 → 30% ⇄ 12 个月退租扣回",
          "链主一事一议最高 150% ⇄ 季度评比",
          "末位淘汰 → 下季度佣金档下调 20%"]),
        ("IV · 生态客户", "14–30 天 3 轮", GOLD,
         ["免租 4 → 6 个月 ⇄ 3+3 年合同",
          "装补 300 → 500 ⇄ 不可单方退租",
          "物业费 28 → 26 ⇄ 年涨 3% 不可调",
          "客户类型差异化报价（算法/硬件/早期）",
          "踏勘后 3–5 天内主动让一步"]),
    ]
    cw = Inches(4.10); cy = Inches(1.20); ch = Inches(5.5); gx = Inches(0.10)
    for i, (t, sub, c, bullets) in enumerate(cards):
        x = Inches(0.5) + (cw + gx) * i
        add_rect(s, x, cy, cw, ch, fill=CLOUD, line=LINE)
        add_rect(s, x, cy, cw, Inches(1.05), fill=c)
        add_text(s, x, cy + Inches(0.10), cw, Inches(0.5), t,
                 size=18, bold=True, color=GOLD if c == NAVY else WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, cy + Inches(0.55), cw, Inches(0.4), sub,
                 size=12, color=GOLD if c == NAVY else WHITE, align=PP_ALIGN.CENTER)
        for j, b in enumerate(bullets):
            yy = cy + Inches(1.30) + Inches(0.85) * j
            add_round(s, x + Inches(0.20), yy + Inches(0.18), Inches(0.30), Inches(0.30),
                      str(j + 1), fill=c, color=WHITE if c != GOLD else NAVY,
                      size=12, bold=True)
            add_text(s, x + Inches(0.55), yy, cw - Inches(0.7), Inches(0.85),
                     b, size=11, color=INK)

    # ============ 新 5. 冠松集团资源协同地图 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="附 · 资源",
               page_title="冠松集团资源协同地图",
               subtitle="把集团 4S 网络 / 保险 / 二手车 / 融资租赁 转化为园区独家壁垒")
    # 中心：冠松
    cx = Inches(6.0); cy = Inches(3.6)
    add_rect(s, cx, cy, Inches(2.4), Inches(0.8), fill=NAVY)
    add_text(s, cx, cy, Inches(2.4), Inches(0.8), "冠松集团",
             size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 6 个外围资源
    resources = [
        ("4S 经销网络", "华东 60+ 网点 · 鸿蒙智行 / 主流品牌\n→ 链主售后渠道协同", BLUE,
         Inches(0.5), Inches(1.2)),
        ("保险事业部", "智驾保险 / 定损 / 理赔大数据\n→ Apollo / Momenta 数据闭环",  PURPLE,
         Inches(9.5), Inches(1.2)),
        ("二手智驾车", "二手车检测 / 翻新 / 流通\n→ 智驾车残值数据 + 数据采集",  GOLD,
         Inches(0.5), Inches(3.3)),
        ("融资租赁子公司", "测试车队融资租赁 / 设备分期\n→ 链主测试车队成本压缩",  GREEN,
         Inches(9.5), Inches(3.3)),
        ("冠松车队（试运营）", "现役车队 100+ · 数据采集与标注\n→ 城市 NOA 数据池",  BLUE,
         Inches(0.5), Inches(5.4)),
        ("冠松产业基金", "战投部 + 5 亿规模\n→ 园区 LP/GP 双模式 + 跟投权", NAVY,
         Inches(9.5), Inches(5.4)),
    ]
    for name, desc, c, x, y in resources:
        add_rect(s, x, y, Inches(3.3), Inches(1.4), fill=CLOUD, line=LINE)
        add_rect(s, x, y, Inches(3.3), Inches(0.45), fill=c)
        add_text(s, x, y, Inches(3.3), Inches(0.45), name,
                 size=12, bold=True, color=NAVY if c == GOLD else WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.1), y + Inches(0.50), Inches(3.1), Inches(0.85),
                 desc, size=10, color=INK)
        # 连线到中心
        line_shp = s.shapes.add_connector(1, x + Inches(1.65), y + Inches(0.7),
                                          cx + Inches(1.2), cy + Inches(0.4))
        line_shp.line.color.rgb = GOLD
        line_shp.line.width = Pt(1.5)

    # ============ 新 6. 22 人组织架构图 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="附 · 团队",
               page_title="22 人稳态组织架构图",
               subtitle="启动 4 人 → 扩张 14 人 → 稳态 22 人 · 8 个部门")
    # 集团董事长
    cx = Inches(5.5); cw = Inches(2.5)
    add_round(s, cx, Inches(1.05), cw, Inches(0.45),
              "集团董事长", fill=NAVY, color=GOLD, size=14, bold=True)
    # 项目总监
    add_round(s, cx, Inches(1.75), cw, Inches(0.45),
              "项目总监 GM (1)", fill=GOLD, color=NAVY, size=14, bold=True)
    # 连线
    line_shp = s.shapes.add_connector(1, cx + Inches(1.25), Inches(1.5),
                                      cx + Inches(1.25), Inches(1.75))
    line_shp.line.color.rgb = NAVY
    line_shp.line.width = Pt(1.5)

    # 8 个部门
    depts = [
        ("GR / 政府事务", "GR 总监\n+ 政府事务经理", "2 人", BLUE,    Inches(0.5),  Inches(2.6)),
        ("招商部", "招商总监\n+ 高招(链主)\n+ 招商×4 + 销售助理\n+ 客户成功", "8 人", GOLD,    Inches(2.05), Inches(2.6)),
        ("运营部", "运营总监\n+ 物业 / IT / 测试场\n+ 行政主任", "5 人", PURPLE,  Inches(3.6),  Inches(2.6)),
        ("市场 / 品牌", "市场总监\n+ 品牌 / 活动经理", "3 人", GREEN,  Inches(5.15), Inches(2.6)),
        ("法务 / 合规", "法务经理", "1 人", RED,    Inches(6.7),  Inches(2.6)),
        ("财务", "财务经理", "1 人", BLUE,   Inches(8.25), Inches(2.6)),
        ("HR / 行政", "HR 经理", "1 人", NAVY,   Inches(9.8),  Inches(2.6)),
        ("项目办公室", "项目总监", "1 人", GOLD,   Inches(11.35),Inches(2.6)),
    ]
    for name, roles, headcount, c, x, y in depts:
        # 部门名
        add_rect(s, x, y, Inches(1.5), Inches(0.45), fill=c)
        add_text(s, x, y, Inches(1.5), Inches(0.45), name,
                 size=11, bold=True, color=NAVY if c == GOLD else WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 人数
        add_round(s, x + Inches(0.4), y + Inches(0.50), Inches(0.7), Inches(0.30),
                  headcount, fill=NAVY, color=GOLD, size=10, bold=True)
        # 角色
        add_rect(s, x, y + Inches(0.85), Inches(1.5), Inches(2.5), fill=CLOUD, line=LINE)
        add_text(s, x + Inches(0.05), y + Inches(0.92), Inches(1.4), Inches(2.4),
                 roles, size=9, color=INK)
        # 连线到项目总监
        line_shp = s.shapes.add_connector(1, x + Inches(0.75), y,
                                          cx + Inches(1.25), Inches(2.20))
        line_shp.line.color.rgb = LINE
        line_shp.line.width = Pt(0.75)

    # 底部统计
    add_rect(s, Inches(0.5), Inches(6.40), Inches(12.3), Inches(0.55), fill=NAVY)
    add_text(s, Inches(0.7), Inches(6.42), Inches(12.0), Inches(0.50),
             "扩编节奏：M1–3 启动 4 人 → M4–6 扩 14 人 → M7–12 稳态 22 人 · 年度人力预算 ¥400→¥1,500→¥2,200 万",
             size=11, bold=True, color=GOLD, anchor=MSO_ANCHOR.MIDDLE)

    # ============ 新 7. 22 人薪酬带宽 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="附 · 团队",
               page_title="22 人岗位薪酬带宽（年度税前 · 万元）",
               subtitle="按等级分组 · D 总监级 / S 资深 / M 经理 / E 普通")
    sal_header = ["#", "岗位", "等级", "年包带宽", "Y3 中位包", "入职时点"]
    sal_rows = [
        ["1", "项目总监 GM", "D", "120–180", "150", "M1"],
        ["2", "GR 总监", "D", "80–130", "105", "M1"],
        ["3", "招商总监", "D", "80–120", "100", "M1"],
        ["4", "运营总监", "D", "60–90", "75", "M1"],
        ["5", "市场总监", "D", "60–85", "72", "M4"],
        ["6", "高级招商经理（链主）", "S", "50–80", "65", "M4"],
        ["7", "法务经理", "S", "45–70", "57", "M4"],
        ["8", "IT / 数据经理", "S", "40–60", "50", "M4"],
        ["9", "测试场协同主任", "S", "40–55", "47", "M4"],
        ["10", "财务经理", "S", "40–60", "50", "M4"],
        ["11", "物业主管", "S", "25–35", "30", "M4"],
        ["12", "招商经理 ×4 (算法/硬件/中介/自拓)", "M", "30–45", "38", "M4"],
        ["13", "政府事务经理", "M", "40–60", "50", "M4"],
        ["14", "客户成功经理", "M", "35–50", "42", "M7"],
        ["15", "品牌 / 活动经理 ×2", "M", "25–40", "32", "M4/M7"],
        ["16", "行政主任 / HR 经理", "M", "20–45", "32", "M7"],
        ["17", "销售助理 / CRM", "E", "15–22", "19", "M4"],
    ]
    add_table(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.0),
              sal_header, sal_rows,
              col_widths=[Inches(0.5), Inches(4.5), Inches(0.8), Inches(2.0),
                          Inches(1.5), Inches(3.0)],
              header_size=11, body_size=10)
    # 等级图例
    add_rect(s, Inches(0.5), Inches(6.30), Inches(12.3), Inches(0.55), fill=NAVY)
    add_text(s, Inches(0.7), Inches(6.32), Inches(12.0), Inches(0.50),
             "★ 链主签约项目奖：项目总监 0.5–0.8% · 招商总监 0.3–0.5% · GR 总监 0.2–0.3% · 期权池核心 8 人 5–8% / 4 年 vesting",
             size=11, bold=True, color=GOLD, anchor=MSO_ANCHOR.MIDDLE)

    # ============ 新 8. KPI 仪表盘 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="附 · KPI",
               page_title="12 个月 KPI 仪表盘",
               subtitle="周度更新 · 月度复盘 · 季度董事会校准")
    kpis_dash = [
        ("团队人数", "M3 4 / M6 14 / M9 18 / M12 22", BLUE),
        ("链主 Term Sheet", "M3 1 / M6 2 / M9 3 / M12 4", GOLD),
        ("链主签约", "M3 0 / M6 1 / M9 2 / M12 3", NAVY),
        ("生态签约 (累计)", "M3 0 / M6 12 / M9 30 / M12 60", PURPLE),
        ("入驻率", "M6 18% / M9 45% / M12 65%", GREEN),
        ("漏斗线索 (累计)", "M3 100 / M6 250 / M9 350 / M12 500", BLUE),
        ("政府汇报次数", "M3 4 / M6 8 / M9 12 / M12 16", GOLD),
        ("媒体曝光指数", "M3 1x / M6 3x / M9 8x / M12 12x", PURPLE),
        ("客户 NPS", "M6 60 / M9 65 / M12 70", GREEN),
        ("现金安全垫 (月)", "M3 12 / M6 9 / M9 7 / M12 8", NAVY),
        ("收入 (累计 · 万)", "M6 250 / M9 600 / M12 1,022", BLUE),
        ("EBITDA 率", "M6 −80% / M9 −60% / M12 −206%", RED),
    ]
    cols = 4; rows = 3
    cw = Inches(3.10); ch = Inches(1.55); gx = Inches(0.05); gy = Inches(0.10)
    for idx, (lbl, val, c) in enumerate(kpis_dash):
        col = idx % cols; row = idx // cols
        x = Inches(0.5) + (cw + gx) * col
        y = Inches(1.20) + (ch + gy) * row
        add_rect(s, x, y, cw, ch, fill=CLOUD, line=LINE)
        add_rect(s, x, y, cw, Inches(0.40), fill=c)
        add_text(s, x, y, cw, Inches(0.40), lbl,
                 size=12, bold=True, color=NAVY if c == GOLD else WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.15), y + Inches(0.50), cw - Inches(0.3), Inches(1.0),
                 val, size=11, color=INK)

    add_rect(s, Inches(0.5), Inches(6.50), Inches(12.3), Inches(0.45), fill=NAVY)
    add_text(s, Inches(0.7), Inches(6.50), Inches(12.0), Inches(0.45),
             "周一招商例会 · 周三 GR 例会 · 周五项目复盘 · 月度运营会 · 季度董事会",
             size=11, bold=True, color=GOLD, anchor=MSO_ANCHOR.MIDDLE)

    # ============ 新 9. Q&A · 高频问答 FAQ ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="附 · FAQ",
               page_title="高频问答 FAQ · 投决会 / 政府汇报常见 10 问",
               subtitle="提前备答 · 有数据 / 有口径 / 有书面材料兜底")
    faq_header = ["#", "常见问题", "应对口径"]
    faq_rows = [
        ["1", "C6 用地下能注册链主总部吗？", "已与区规划/经委预先沟通\"用途相符性\"，链主主体名单经备案"],
        ["2", "Y3 入驻率 92% 会做不到吗？", "敏感性：入驻 -10% 仍能 Y3 EBITDA 转正；最坏情形不会黑洞"],
        ["3", "链主 TOP5 跑了怎么办？", "备选名单（蔚来/理想/大众问问/Smart）+ 8/9F 分租方案"],
        ["4", "户外测试场没有，能算智驾园吗？", "三段式方案：园区静态 + 1.5km 路测延伸 + 嘉定/临港会员"],
        ["5", "政策返还万一兑现不了？", "区财政书面确认函；改分期口径；不构成项目方独立担保"],
        ["6", "9 月发布会能拉多少政府嘉宾？", "目标 25 位（市级 1 + 区级 5 + 委办 19），区委书记/区长背书"],
        ["7", "Y0 启动现金需求多少？", "约 3 亿；Y1 期初安全垫 ≥ 12 个月；Y3 期末现金 + 1.28 亿"],
        ["8", "如果价格被压到 5.0 元/㎡·天？", "已超底线红线；立即休会请示集团董事长"],
        ["9", "公关与媒体年度预算多少？", "1,500–2,000 万/年；目标全媒体阅读 ≥ 1 亿，央/新华 ≥ 6 篇"],
        ["10", "退出机制有吗？", "Y3 后视入驻率推 REITs 化；集团信用借款托底；A 栋分租预案"],
    ]
    add_table(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5),
              faq_header, faq_rows,
              col_widths=[Inches(0.5), Inches(4.5), Inches(7.3)],
              header_size=11, body_size=10)
    add_rect(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.10), fill=GOLD)

    # ============ 新 10. 文档索引 · 决策包 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="附 · 索引",
               page_title="决策文档包 · 一键直达",
               subtitle="本汇报后续支撑材料清单 · 法务 / 财务 / 商务 / 团队 全覆盖")

    add_text(s, Inches(0.5), Inches(1.10), Inches(6), Inches(0.4),
             "📑 商务汇报 PPT 与策略", size=14, bold=True, color=NAVY)
    pkgs = [
        ("项目总览", "README.md", BLUE),
        ("产业定位报告", "docs/phase1-strategy/01-industry-positioning-report.md", BLUE),
        ("空间规划（9F 楼层）", "docs/phase1-strategy/02-space-planning.md", BLUE),
        ("链主攻坚 + 5 张一页纸", "docs/phase2-execution/03-anchor-tenant-tracker.md", GOLD),
        ("谈判策略 Playbook", "docs/phase2-execution/03b-negotiation-playbook.md", GOLD),
        ("生态漏斗 + 中介策略", "docs/phase2-execution/04-ecosystem-funnel.md", GOLD),
        ("政府关系 + 政策包", "docs/phase2-execution/05-government-relations.md", GOLD),
        ("品牌活动 + 9 月发布会", "docs/phase3-brand/06-launch-and-events.md", PURPLE),
        ("商业条款 + 报价单", "docs/phase4-commercial/07-pricing-and-contract.md", GREEN),
        ("12 个月执行计划", "docs/phase5-rollout/08-execution-plan.md", NAVY),
        ("22 人 JD + 薪酬 + KPI", "docs/phase5-rollout/08b-team-and-jd.md", NAVY),
    ]
    for i, (name, path, c) in enumerate(pkgs):
        y = Inches(1.55) + Inches(0.42) * i
        add_rect(s, Inches(0.5), y, Inches(0.18), Inches(0.34), fill=c)
        add_rect(s, Inches(0.7), y, Inches(6.0), Inches(0.34), fill=CLOUD, line=LINE)
        add_text(s, Inches(0.85), y, Inches(2.5), Inches(0.34), name,
                 size=10, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(3.30), y, Inches(3.30), Inches(0.34), path,
                 size=8, color=GREY, anchor=MSO_ANCHOR.MIDDLE, italic=True)

    add_text(s, Inches(7.0), Inches(1.10), Inches(6), Inches(0.4),
             "📋 法务 Word + 财务 Excel", size=14, bold=True, color=NAVY)
    legal_pkgs = [
        ("链主总部租赁合同 (草案)", "docs/legal/01-合作协议-链主总部租赁合同.docx", RED),
        ("中介居间服务协议 (草案)", "docs/legal/02-合作协议-中介居间服务协议.docx", RED),
        ("联合实验室共建协议 (草案)", "docs/legal/03-合作协议-联合实验室共建协议.docx", RED),
        ("政府专班合作备忘录 (草案)", "docs/legal/04-合作协议-政府专班合作备忘录.docx", RED),
        ("财务测算与商务模型 (8 Sheet)", "docs/finance/财务测算与商务模型.xlsx", GOLD),
        ("─ Sheet 1：摘要 Dashboard", "—", GREY),
        ("─ Sheet 2：假设与参数", "—", GREY),
        ("─ Sheet 3：三年损益", "—", GREY),
        ("─ Sheet 4：月度滚动现金流（36m）", "—", GREY),
        ("─ Sheet 5：入驻进度爬坡（36m）", "—", GREY),
        ("─ Sheet 6：双变量敏感性", "—", GREY),
        ("─ Sheet 7：22 人薪酬带宽", "—", GREY),
        ("─ Sheet 8：让步阶梯计算器", "—", GREY),
    ]
    for i, (name, path, c) in enumerate(legal_pkgs):
        y = Inches(1.55) + Inches(0.36) * i
        add_rect(s, Inches(7.0), y, Inches(0.18), Inches(0.30), fill=c)
        add_rect(s, Inches(7.2), y, Inches(5.7), Inches(0.30), fill=CLOUD, line=LINE)
        add_text(s, Inches(7.35), y, Inches(2.6), Inches(0.30), name,
                 size=9, bold=(c != GREY), color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(9.95), y, Inches(2.95), Inches(0.30), path,
                 size=8, color=GREY, anchor=MSO_ANCHOR.MIDDLE, italic=True)

    # ============ 34. 投决建议 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="07 · 投决建议",
               page_title="投决建议 · 里程碑承诺",
               subtitle="资源到位的前提下，T+30 / 90 / 180 三个节点的硬承诺")
    nodes2 = [
        ("T + 30 天", "团队就位 + 政府首轮汇报",
         ["4 人核心入职", "区投促办 + 副区长 + 市经信委首轮汇报", "5 家中介签约"], BLUE),
        ("T + 90 天", "链主首份 Term Sheet",
         ["TOP5 全部进入会面阶段", "≥ 1 家链主签 Term Sheet", "300 家库 + CRM 上线"], GOLD),
        ("T + 180 天", "9 月发布会 + 政策包定稿",
         ["发布会 200 人到场 + 5 家签约", "「iDrive 静安 10 条」发布", "签约 ≥ 30 家 / 入驻率 ≥ 30%"], GREEN),
    ]
    for i, (t, sub, bullets, c) in enumerate(nodes2):
        x = Inches(0.5) + Inches(4.30) * i
        add_rect(s, x, Inches(1.2), Inches(4.10), Inches(4.5), fill=CLOUD, line=LINE)
        add_rect(s, x, Inches(1.2), Inches(4.10), Inches(0.95), fill=c)
        add_text(s, x, Inches(1.25), Inches(4.10), Inches(0.5), t,
                 size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(1.70), Inches(4.10), Inches(0.4), sub,
                 size=12, color=WHITE, align=PP_ALIGN.CENTER)
        for j, b in enumerate(bullets):
            y = Inches(2.35) + Inches(0.65) * j
            add_round(s, x + Inches(0.20), y + Inches(0.15), Inches(0.25), Inches(0.25),
                      str(j + 1), fill=c, color=WHITE, size=10, bold=True)
            add_text(s, x + Inches(0.55), y + Inches(0.05), Inches(3.40), Inches(0.55),
                     b, size=12, color=INK)
    add_rect(s, Inches(0.5), Inches(5.95), Inches(12.3), Inches(1.0), fill=NAVY)
    add_text(s, Inches(0.7), Inches(6.0), Inches(12.0), Inches(0.45),
             "请董事会审议", size=14, bold=True, color=GOLD)
    add_text(s, Inches(0.7), Inches(6.40), Inches(12.0), Inches(0.55),
             "① 启动预算（Y0–Y1）  ② 4 人核心团队招聘授权  ③ 静安区一企一策政府专班 ④ 9 月发布会预算  ⑤ 链主谈判授权区间",
             size=12, color=WHITE)

    # ============ 35. Q&A ============
    s = new_slide(prs)
    add_rect(s, 0, 0, sw, sh, fill=NAVY)
    add_rect(s, 0, Inches(3.6), sw, Emu(40000), fill=GOLD)
    add_text(s, Inches(0.8), Inches(2.3), Inches(11), Inches(1.4),
             "GS · iDrive Hub", size=44, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), Inches(3.8), Inches(11), Inches(0.7),
             "让中心城区跑通智能驾驶最后一公里",
             size=22, color=WHITE, italic=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), Inches(5.0), Inches(11), Inches(0.5),
             "Q & A · 谢 谢", size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), Inches(6.4), Inches(11), Inches(0.4),
             "冠松集团 · GS iDrive Hub 项目组",
             size=12, color=CLOUD, align=PP_ALIGN.CENTER)

    # --------------------------------------------------------------------
    # 回填总页数
    # --------------------------------------------------------------------
    total = len(SLIDES)
    for sl in SLIDES:
        for shape in sl.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text and run.text.strip().endswith(" / 0"):
                        run.text = run.text.replace(" / 0", f" / {total}")

    out = Path(__file__).resolve().parent.parent / "docs" / "deck" / \
        "GS-iDrive-Hub-招商方案.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"✓ Deck written: {out}  ({total} slides)")
    return out


if __name__ == "__main__":
    build()
