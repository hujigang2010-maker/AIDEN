"""生成 PPT：AI+机器人轻资产创业活动策划案（紫色主题）"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# 紫色主题配色
PRIMARY = RGBColor(0x4A, 0x14, 0x8C)       # 深紫
ACCENT = RGBColor(0x7B, 0x1F, 0xA2)        # 中紫
LIGHT = RGBColor(0xF3, 0xE5, 0xF5)         # 浅紫底
LAVENDER = RGBColor(0xCE, 0x93, 0xD8)      # 薰衣草紫
GOLD = RGBColor(0xE1, 0xBE, 0xE7)          # 淡紫高光
DARK = RGBColor(0x2D, 0x1B, 0x3D)            # 深紫文字
GREY = RGBColor(0x6A, 0x5A, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_DEEP = RGBColor(0x31, 0x0A, 0x5C)        # 封面背景


def set_font(run, name="微软雅黑", size=18, bold=False, color=DARK):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    rFonts = rPr.find(qn("a:ea"))
    if rFonts is None:
        rFonts = rPr.makeelement(qn("a:ea"), {"typeface": name})
        rPr.append(rFonts)
    else:
        rFonts.set("typeface", name)


def add_rect(slide, left, top, width, height, fill=PRIMARY, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, size=18, bold=False,
             color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font="微软雅黑"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        text = [text]
    for i, t in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = t
        set_font(run, name=font, size=size, bold=bold, color=color)
    return tb


def add_bullet_list(slide, left, top, width, height, items, size=14,
                    color=DARK, bullet="•", line_spacing=1.25):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = f"{bullet}  {it}"
        set_font(run, size=size, color=color)
    return tb


def slide_header(slide, title, subtitle=None, page_no=None, total=None):
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.45), fill=PRIMARY)
    add_rect(slide, Inches(0), Inches(0.45), Inches(13.333), Inches(0.05), fill=LAVENDER)
    add_text(slide, Inches(0.5), Inches(0.55), Inches(10), Inches(0.6),
             title, size=24, bold=True, color=PRIMARY, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.05), Inches(12), Inches(0.35),
                 subtitle, size=12, color=GREY, anchor=MSO_ANCHOR.MIDDLE)
    if page_no and total:
        add_text(slide, Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.3),
                 f"{page_no} / {total}", size=10, color=GREY, align=PP_ALIGN.RIGHT)


def make_card(slide, left, top, width, height, title, body_items,
              accent=ACCENT, title_size=14, body_size=11):
    add_rect(slide, left, top, width, height, fill=LIGHT)
    add_rect(slide, left, top, Inches(0.08), height, fill=accent)
    add_text(slide, left + Inches(0.2), top + Inches(0.1),
             width - Inches(0.3), Inches(0.4),
             title, size=title_size, bold=True, color=accent)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.55),
                    width - Inches(0.3), height - Inches(0.6),
                    body_items, size=body_size, color=DARK)


def make_table(slide, left, top, width, height, headers, rows, header_color=PRIMARY, font_size=10):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    table = shape.table
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.05)
        cell.margin_bottom = Inches(0.05)
        tf = cell.text_frame
        tf.text = ""
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = h
        set_font(r, size=font_size + 1, bold=True, color=WHITE)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 == 0 else LIGHT
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            tf = cell.text_frame
            tf.text = ""
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(val)
            set_font(r, size=font_size, color=DARK)
    return table


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    TOTAL = 18
    page = [0]

    def new_slide():
        page[0] += 1
        return prs.slides.add_slide(blank), page[0]

    # 1 封面
    s, _ = new_slide()
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), fill=BG_DEEP)
    add_rect(s, Inches(0), Inches(3.5), Inches(13.333), Inches(0.06), fill=LAVENDER)
    add_rect(s, Inches(0.8), Inches(1.2), Inches(0.12), Inches(2.2), fill=ACCENT)
    add_text(s, Inches(1.1), Inches(1.5), Inches(11.5), Inches(1.0),
             "AI + 机器人轻资产创业活动", size=48, bold=True, color=WHITE)
    add_text(s, Inches(1.1), Inches(2.6), Inches(11.5), Inches(0.8),
             "寒假研学科创营 · 策划案", size=36, bold=True, color=GOLD)
    add_text(s, Inches(1.1), Inches(3.8), Inches(11.5), Inches(0.5),
             "资源池 × 产品 × 客群对齐  |  轻资产验证  |  研学场景切入",
             size=16, color=LAVENDER)
    add_text(s, Inches(1.1), Inches(6.0), Inches(11.5), Inches(0.4),
             "青岛西海岸（黄岛）试点  |  V1.0 讨论稿  |  2025年7月",
             size=12, color=WHITE)
    add_text(s, Inches(1.1), Inches(6.5), Inches(11.5), Inches(0.4),
             "基于《AI与机器人创业方向探讨》会议纪要整理",
             size=11, color=GREY)

    # 2 目录
    s, p = new_slide()
    slide_header(s, "目录", "讨论用策划案结构", p, TOTAL)
    items = [
        "一、会议纪要梗概",
        "二、创业方向与战略判断",
        "三、活动定位与目标客群",
        "四、商业模式与产品组合",
        "五、合作方策略（高校 / 黄岛区科技企业）",
        "六、七至十二月倒排计划",
        "七、寒假主活动方案",
        "八、预算与风险控制",
        "九、待讨论事项与下一步",
    ]
    add_bullet_list(s, Inches(1.5), Inches(1.8), Inches(10), Inches(5),
                    items, size=20, color=DARK, bullet="▸")

    # 3 会议纪要梗概
    s, p = new_slide()
    slide_header(s, "一、会议纪要梗概", "两份材料核心要点提炼", p, TOTAL)
    make_card(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.2),
              "会议背景",
              ["地产转型团队探讨进入 AI 与机器人赛道",
               "现有资源：研学渠道（学校/教育机构）、房企上下游私企",
               "原计划：机器人租售 + AI 软件双线、混改国企建展厅",
               "对接品牌：深蓝科技（机器人）、智巨人（AI 视频）",
               "区域：青岛黄岛区，目标寒假落地主活动"])
    make_card(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.2),
              "专家核心建议",
              ["地产不宜再作主业，须转向科技轻资产",
               "暂缓展厅与混改，优先走出去办活动",
               "资源池·产品·客群三要素必须对齐",
               "机器人与 AI 客群不同，避免双线分散",
               "优选头部品牌，一级代理+提成，勿压库存",
               "研学是可控资源，但需拓展 B 端变现场景"],
              accent=PRIMARY)

    # 4 战略判断
    s, p = new_slide()
    slide_header(s, "二、创业方向与战略判断", "从地产思维转向需求驱动", p, TOTAL)
    make_card(s, Inches(0.5), Inches(1.5), Inches(3.9), Inches(2.3),
              "行业判断",
              ["地产长期横盘，仅作辅助变现",
               "科技驱动型赛道处于红利期",
               "马太效应加速，须绑定头部厂商"])
    make_card(s, Inches(4.7), Inches(1.5), Inches(3.9), Inches(2.3),
              "资源判断",
              ["政府/国企资源不可控，易被架空",
               "研学渠道是团队最大可控资产",
               "B 端（婚庆/活动/短视频）可延伸"])
    make_card(s, Inches(8.9), Inches(1.5), Inches(3.9), Inches(2.3),
              "模式判断",
              ["游击战 > 阵地战（无固定展厅）",
               "PPT + 视频 + 样机演示即可",
               "先验证 MVP，再考虑重投入"])
    make_card(s, Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.6),
              "推荐主路径（讨论稿）",
              ["聚焦「研学 + AI 科创体验」统一客群：中小学生研学团 + 教育机构",
               "机器人以租赁/联合演示为主，不自购库存；AI 以短视频创作体验课引流算力续费",
               "参考青岛「萝卜快跑」模式：小工作室 + 免费体验课 + 流量包续费",
               "暂缓 50 万级区域代理买断，优先谈一级代理提成或活动分成"])

    # 5 活动定位
    s, p = new_slide()
    slide_header(s, "三、活动定位与目标客群", "寒假主活动锚点", p, TOTAL)
    add_text(s, Inches(0.5), Inches(1.5), Inches(12), Inches(0.5),
             "活动名称（建议）：「智创未来」寒假 AI+机器人科创研学营",
             size=20, bold=True, color=PRIMARY)
    make_table(s, Inches(0.5), Inches(2.2), Inches(12.3), Inches(3.8),
               ["客群", "需求场景", "产品/服务", "变现方式"],
               [["研学学生（7-15岁）", "寒假科创打卡、机器人体验", "半日营/一日营", "B端机构采购+门票分成"],
                ["教育机构/学校", "丰富研学产品线", "定制研学包", "按人头服务费"],
                ["活动公司/商场", "引流表演、节庆活动", "机器人租赁+主持", "单场租赁费"],
                ["小微创业团队", "AI短视频制作", "体验课+算力包", "课程费+流量包续费"],
                ["家长社群", "科技启蒙、升学背景", "亲子工作坊", "报名费+衍生品"]],
               header_color=ACCENT)

    # 6 商业模式
    s, p = new_slide()
    slide_header(s, "四、商业模式与产品组合", "轻资产三层收入结构", p, TOTAL)
    make_card(s, Inches(0.5), Inches(1.5), Inches(3.9), Inches(4.8),
              "第一层：活动服务费",
              ["寒假研学营学费/机构采购",
               "商场/节庆单场活动执行费",
               "企业 Open Day 定制体验"])
    make_card(s, Inches(4.7), Inches(1.5), Inches(3.9), Inches(4.8),
              "第二层：代理提成",
              ["机器人租赁/销售分成（宇树/智谱等）",
               "AI 算力/即梦等流量包分销",
               "不压货，成交后结算"])
    make_card(s, Inches(8.9), Inches(1.5), Inches(3.9), Inches(4.8),
              "第三层：生态合作",
              ["教育机构长期课程合作",
               "科技企业联合品牌赞助",
               "高校科普基地挂牌分成"])
    add_text(s, Inches(0.5), Inches(6.5), Inches(12), Inches(0.5),
             "原则：同一客群可「研学体验 + AI 短课」组合；机器人与 AI 算力分渠道运营，避免双线抢精力",
             size=12, color=GREY)

    # 7 合作方策略
    s, p = new_slide()
    slide_header(s, "五、合作方策略", "高校中心与黄岛区科技企业 — 建议方案", p, TOTAL)
    make_card(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.0),
              "✅ 建议引入：高校/中航科幻类中心",
              ["角色：科普内容背书 + 研学课程共建 + 师资/志愿者",
               "价值：提升活动公信力，满足学校研学合规需求",
               "方式：联合挂牌「科创研学实践基地」、提供科普讲座",
               "节奏：9月洽谈意向，10月签约框架，11月联合宣发",
               "注意：不参与股权与重资产投入，以内容合作为主"],
              accent=ACCENT)
    make_card(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.0),
              "✅ 建议引入：黄岛区科技企业",
              ["角色：样机支持、技术讲解、品牌联合露出",
               "价值：解决机器人演示供应链，降低自购成本",
               "方式：加入区人工智能/机器人协会，横向比选合作方",
               "节奏：8月协会对接，9月确定 2-3 家联合演示伙伴",
               "备选：扬州等外地科技企业可作为设备租赁补充"],
              accent=PRIMARY)
    add_text(s, Inches(0.5), Inches(6.7), Inches(12), Inches(0.4),
             "不建议：以政府/国企混改为主体推动；可作为场地赞助方，不做核心依赖",
             size=12, bold=True, color=ACCENT)

    # 8 倒排计划总览
    s, p = new_slide()
    slide_header(s, "六、七至十二月倒排计划", "以 2026 寒假为主节点逆向拆解", p, TOTAL)
    rows = [
        ["12月", "冲刺筹备", "场地确认、师资到位、全流程彩排、寒假招生截止"],
        ["11月", "招生预热", "研学机构签约、海报投放、体验课开放日、物料到位"],
        ["10月", "产品封装", "课程大纲定稿、PPT/视频物料、报价体系、合同模板"],
        ["9月", "试点验证", "小规模体验场（2-3场）、合作方签约、代理政策落地"],
        ["8月", "资源对接", "协会入会、高校/企业洽谈、品牌比选、团队分工"],
        ["7月", "战略对齐", "本策划案讨论定稿、调研头部品牌、研学资源盘点"],
    ]
    make_table(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(4.5),
               ["月份", "阶段主题", "关键里程碑"],
               rows, header_color=PRIMARY, font_size=11)
    add_text(s, Inches(0.5), Inches(6.3), Inches(12), Inches(0.5),
             "▶ 主活动执行窗口：2026年1月中下旬—2月（寒假），详见下页",
             size=14, bold=True, color=ACCENT)

    # 9 各月详细节点
    s, p = new_slide()
    slide_header(s, "倒排计划 · 各月关键任务", "责任分工建议在 Excel 中细化", p, TOTAL)
    make_table(s, Inches(0.3), Inches(1.4), Inches(12.7), Inches(5.5),
               ["月份", "核心任务", "交付物", "决策点"],
               [
                   ["7月", "战略对齐、品牌调研、资源盘点", "策划案定稿、品牌短名单", "是否暂缓展厅/混改"],
                   ["8月", "协会对接、高校/企业拜访", "合作意向书、分工表", "选定 2-3 家演示伙伴"],
                   ["9月", "MVP 体验活动 2-3 场", "试点复盘、代理条款", "确定主代理品牌"],
                   ["10月", "课程产品化、物料制作", "课程包、宣传册、报价单", "寒假产品定价"],
                   ["11月", "招生与渠道签约", "机构合同、报名数据", "是否达到最低开班人数"],
                   ["12月", "彩排与运营准备", "执行手册、应急预案", "寒假档期最终确认"],
               ],
               header_color=ACCENT, font_size=9)

    # 10 寒假主活动
    s, p = new_slide()
    slide_header(s, "七、寒假主活动方案", "2026 寒假 · 核心交付", p, TOTAL)
    make_card(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(2.5),
              "活动形态",
              ["3-5 期寒假科创研学营（半日/一日）",
               "每期 30-50 人，合计目标 150-250 人次",
               "模块：机器人互动 + AI 短视频体验 + 科普讲座"])
    make_card(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.5),
              "场地策略",
              ["优先：合作商场中庭/社区活动中心/高校科普场地",
               "备选：黄岛区国企闲置商业（仅租场地，不建展厅）",
               "原则：游击式多点位，不集中单一展厅"])
    make_card(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.5),
              "执行流程（单日营示例）",
              ["09:00-09:30 签到破冰  →  09:30-10:30 机器人表演与互动体验",
               "10:45-11:45 AI 短视频创作小课堂（即梦/类似工具）",
               "13:30-15:00 分组科创挑战 + 成果展示  →  15:00-15:30 结业颁证/合影"])

    # 11 品牌与供应链
    s, p = new_slide()
    slide_header(s, "品牌与供应链建议", "基于会议纪要的品牌研判", p, TOTAL)
    make_table(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(4.0),
               ["品类", "推荐方向", "暂不建议", "合作模式"],
               [
                   ["人形/表演机器人", "宇树、智谱等头部", "深蓝 50 万买断加盟", "租赁+销售提成"],
                   ["AI 视频/算力", "字节即梦、DeepSeek 等", "智巨人等非头部", "一级代理 4-5 折"],
                   ["行业资源", "青岛 AI/机器人协会", "单纯依赖政府关系", "副会长/会员身份"],
               ],
               header_color=PRIMARY)
    add_bullet_list(s, Inches(0.5), Inches(5.8), Inches(12), Inches(1.2),
                    ["库存风险：机器人迭代快（机器狗已四代），严禁大批量囤货",
                     "验证顺序：先租后买，先活动后代理，先研学后 B 端拓展"],
                    size=13, color=DARK)

    # 12 营销打法
    s, p = new_slide()
    slide_header(s, "营销打法：游击战而非阵地战", "走出去，把客户找出来", p, TOTAL)
    make_card(s, Inches(0.5), Inches(1.5), Inches(3.9), Inches(4.5),
              "渠道获客",
              ["研学合作机构批量导入",
               "学校家委会/家长群裂变",
               "商场节庆联合引流"])
    make_card(s, Inches(4.7), Inches(1.5), Inches(3.9), Inches(4.5),
              "内容获客",
              ["短视频探店/活动花絮",
               "小红书/抖音本地生活",
               "孩子作品线上展播"])
    make_card(s, Inches(8.9), Inches(1.5), Inches(3.9), Inches(4.5),
              "活动获客",
              ["周末免费体验课（筛选付费）",
               "B 端开放日（婚庆/活动公司）",
               "科技企业联合路演"])
    add_text(s, Inches(0.5), Inches(6.3), Inches(12), Inches(0.6),
             "关键原则：算力与 AI 产品须走进企业/活动现场推销；线上免费课转化低，须线下收费筛选",
             size=13, bold=True, color=ACCENT)

    # 13 组织架构
    s, p = new_slide()
    slide_header(s, "团队分工建议", "最小可行团队（5-8人）", p, TOTAL)
    make_table(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(4.5),
               ["角色", "职责", "人数", "备注"],
               [
                   ["项目负责人", "战略、合作谈判、预算", "1", "地产转型核心成员"],
                   ["研学渠道", "机构对接、招生、排期", "1-2", "现有核心资源"],
                   ["内容与课程", "课程设计、讲师、物料", "1-2", "可联合高校师资"],
                   ["活动执行", "现场统筹、设备、安全", "2", "可兼职/志愿者"],
                   ["市场运营", "宣发、短视频、数据", "1", "轻资产为主"],
               ],
               header_color=ACCENT)

    # 14 预算
    s, p = new_slide()
    slide_header(s, "八、预算与投入（轻资产版）", "首年验证期 · 不含代理买断", p, TOTAL)
    make_table(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(4.2),
               ["项目", "金额（万元）", "说明", "性质"],
               [
                   ["协会/行业资源", "3-5", "人工智能协会副会长等", "一次性"],
                   ["设备租赁（全年）", "8-15", "机器人按场租赁，非购置", "变动成本"],
                   ["场地与活动", "5-10", "多点位游击式活动", "变动成本"],
                   ["物料与课程开发", "3-5", "PPT、视频、教材", "一次性"],
                   ["市场与招生", "3-5", "宣发、体验课", "变动成本"],
                   ["人员与运营", "10-15", "5-8人兼职/核心团队", "固定成本"],
                   ["预备金", "5", "应急与试点追加", "储备"],
                   ["合计", "37-60", "远低於 50 万单一代理买断", "—"],
               ],
               header_color=PRIMARY, font_size=10)

    # 15 风险
    s, p = new_slide()
    slide_header(s, "风险控制", "会议纪要警示项", p, TOTAL)
    make_table(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(4.8),
               ["风险", "等级", "应对策略"],
               [
                   ["双线作战精力分散", "高", "统一研学客群，分渠道运营"],
                   ["机器人库存贬值", "高", "只租不买，成交后结算"],
                   ["政府资源被架空", "高", "不依赖政府获客，仅作场地赞助"],
                   ["招生不达预期", "中", "11月设最低开班线，未达标延期"],
                   ["品牌方代理政策变化", "中", "多品牌比选，不签独家买断"],
                   ["免费课转化低", "中", "线下收费体验课筛选付费用户"],
               ],
               header_color=ACCENT, font_size=10)

    # 16 KPI
    s, p = new_slide()
    slide_header(s, "核心 KPI（寒假验证期）", "可量化的 MVP 目标", p, TOTAL)
    make_table(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(4.0),
               ["维度", "指标", "目标值", "验证意义"],
               [
                   ["招生", "寒假营总人次", "≥150", "研学渠道变现能力"],
                   ["收入", "活动总收入", "≥15万元", "商业模式可行性"],
                   ["渠道", "签约研学机构", "≥3家", "B端复购基础"],
                   ["合作", "联合企业/高校", "各≥1家", "生态背书"],
                   ["转化", "算力/课程续费用户", "≥30人", "第二曲线验证"],
                   ["成本", "单场活动毛利", "≥25%", "可持续运营"],
               ],
               header_color=PRIMARY)

    # 17 待讨论
    s, p = new_slide()
    slide_header(s, "九、待讨论事项", "供团队评审决策", p, TOTAL)
    add_bullet_list(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.2),
                    [
                        "是否确认以「寒假研学营」为第一个 MVP，暂缓展厅与混改？",
                        "机器人 vs AI 算力：是否同意「统一研学客群、分团队运营」？",
                        "高校/中航科幻中心：是否启动联合研学基地洽谈？",
                        "黄岛区科技企业：是否优先加入区人工智能协会？",
                        "代理品牌：是否将宇树/智谱/即梦列为重点调研对象？",
                        "预算区间 37-60 万是否可接受？是否需要压缩至 30 万以内？",
                        "寒假档期：1月还是2月为主？每期人数与定价？",
                        "是否需要赴上海/行业展会补看品牌（人工智能大会等）？",
                    ],
                    size=16, color=DARK, bullet="?")

    # 18 封底
    s, p = new_slide()
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), fill=BG_DEEP)
    add_rect(s, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06), fill=LAVENDER)
    add_text(s, Inches(1.0), Inches(2.5), Inches(11), Inches(1.0),
             "谢谢 · 欢迎讨论", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.0), Inches(3.8), Inches(11), Inches(0.6),
             "在不熟悉的领域多看看，再决定 —— 人工智能周期还长，不必急于求成",
             size=16, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.0), Inches(5.5), Inches(11), Inches(0.5),
             "附件：Excel 倒排计划表  |  会议纪要完整版",
             size=12, color=LAVENDER, align=PP_ALIGN.CENTER)

    out = "/workspace/output/AI机器人创业活动策划案.pptx"
    import os
    os.makedirs(os.path.dirname(out), exist_ok=True)
    prs.save(out)
    print(f"已生成: {out}")


if __name__ == "__main__":
    main()
