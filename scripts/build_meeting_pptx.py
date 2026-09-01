"""
冠松 01# 楼 · 招商合作会上版 PPT（6 页）

只讲：听懂了、粗思路、收费、怎么停、请定。
不讲 iDrive Hub、不讲链主 TOP5、不讲 22 人团队。

重新生成：python3 scripts/build_meeting_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pathlib import Path

NAVY = RGBColor(0x1B, 0x35, 0x55)
GOLD = RGBColor(0xB7, 0x86, 0x2E)
INK = RGBColor(0x2A, 0x2D, 0x34)
CLOUD = RGBColor(0xF0, 0xF1, 0xF3)
GREY = RGBColor(0x6B, 0x73, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
STONE = RGBColor(0x8A, 0x92, 0x9C)

CN_FONT = "WenQuanYi Micro Hei"
EN_FONT = "Georgia"
SLIDES = []


def set_run(run, text, *, size=14, bold=False, color=INK, italic=False):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = EN_FONT
    rPr = run._r.get_or_add_rPr()
    for tag in ("ea", "cs"):
        existing = rPr.find(qn(f"a:{tag}"))
        if existing is not None:
            rPr.remove(existing)
        rPr.append(rPr.makeelement(qn(f"a:{tag}"), {"typeface": CN_FONT}))


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    box.fill.background()
    box.line.fill.background()
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        set_run(run, ln, size=size, bold=bold, color=color, italic=italic)
    return box


def add_rect(slide, x, y, w, h, *, fill=NAVY):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    shp.text_frame.text = ""
    return shp


def add_round(slide, x, y, w, h, text="", *, fill=NAVY, color=WHITE,
              size=12, bold=True, align=PP_ALIGN.CENTER):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p.add_run(), text, size=size, bold=bold, color=color)
    return shp


def add_table(slide, x, y, w, h, header, rows, *, header_size=12, body_size=12,
              col_widths=None):
    ts = slide.shapes.add_table(len(rows) + 1, len(header), x, y, w, h)
    table = ts.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw
    for j, htxt in enumerate(header):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.text = ""
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p.add_run(), htxt, size=header_size, bold=True, color=WHITE)
    for i, row in enumerate(rows, start=1):
        bg = WHITE if i % 2 else CLOUD
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            cell.text = ""
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            set_run(p.add_run(), str(val), size=body_size, color=INK)
    return ts


def add_chrome(slide, prs, *, page_no, label, title, subtitle=""):
    sw, sh = prs.slide_width, prs.slide_height
    add_rect(slide, 0, 0, sw, Emu(380000), fill=NAVY)
    add_round(slide, Inches(0.45), Inches(0.18), Inches(2.4), Inches(0.36),
              label, fill=GOLD, color=NAVY, size=11)
    add_text(slide, Inches(3.05), Inches(0.10), Inches(9.6), Inches(0.52),
             title, size=22, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(3.05), Inches(0.52), Inches(9.6), Inches(0.28),
                 subtitle, size=12, color=CLOUD, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 0, sh - Emu(80000), sw, Emu(80000), fill=GOLD)
    add_text(slide, Inches(0.45), Inches(7.05), Inches(9.2), Inches(0.28),
             "冠松 01# 楼 · 招商合作会上版 · 保密", size=10, color=GREY)
    add_text(slide, Inches(11.4), Inches(7.05), Inches(1.5), Inches(0.28),
             f"{page_no} / 0", size=10, color=GREY, align=PP_ALIGN.RIGHT)


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    SLIDES.append(s)
    return s


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    sw, sh = prs.slide_width, prs.slide_height

    # 1 封面
    s = new_slide(prs)
    add_rect(s, 0, 0, sw, sh, fill=NAVY)
    add_rect(s, 0, Inches(5.05), sw, Emu(28000), fill=GOLD)
    add_text(s, Inches(0.8), Inches(1.35), Inches(11.5), Inches(0.45),
             "01# 研发楼 · 招商合作", size=16, bold=True, color=GOLD, italic=True)
    add_text(s, Inches(0.8), Inches(1.9), Inches(11.7), Inches(1.4),
             "我们帮你们把这栋楼租出去",
             size=40, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.5), Inches(11.7), Inches(0.7),
             "合作协议讲收费 · 这几页讲思路 · 招到盖章才算完",
             size=20, color=CLOUD)
    add_round(s, Inches(0.8), Inches(5.5), Inches(3.4), Inches(0.48),
              "会上版 · 6 页", fill=GOLD, color=NAVY, size=14)
    add_text(s, Inches(4.4), Inches(5.5), Inches(8), Inches(0.48),
             "配套：05b 商务简版（收费确认栏可当场填）",
             size=14, color=CLOUD, anchor=MSO_ANCHOR.MIDDLE)

    # 2 听懂了
    s = new_slide(prs)
    add_chrome(s, prs, page_no=2, label="对齐",
               title="我们听懂了",
               subtitle="得到大脑口径：跟进招商。不是再做一版定位。")
    items = [
        ("1", "楼要有人租", "汽车展厅走不通了。目标是生效合同，不是品牌发布会。"),
        ("2", "不当小白鼠", "抄已经跑通的研发楼，不搞全国首创。"),
        ("3", "不丢掉汽车", "不招整车 4S；汽车配套、后市场可以留。"),
        ("4", "你们拍板我们跑", "合同您盖章。我们建库、带看、谈判、盯收款。"),
    ]
    for i, (n, t, b) in enumerate(items):
        y = Inches(1.25) + Inches(1.3) * i
        add_round(s, Inches(0.5), y, Inches(0.7), Inches(1.05), n,
                  fill=NAVY, color=WHITE, size=22)
        add_text(s, Inches(1.45), y, Inches(3.3), Inches(1.05), t,
                 size=22, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(4.8), y, Inches(7.9), Inches(1.05), b,
                 size=18, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # 3 粗思路
    s = new_slide(prs)
    add_chrome(s, prs, page_no=3, label="思路",
               title="招谁、怎么招",
               subtitle="C6 教育科研用地是硬约束。地上不能按 4S 店来用。")

    add_rect(s, Inches(0.45), Inches(1.25), Inches(6.1), Inches(5.35), fill=CLOUD)
    add_rect(s, Inches(0.45), Inches(1.25), Inches(6.1), Inches(0.55), fill=NAVY)
    add_text(s, Inches(0.45), Inches(1.25), Inches(6.1), Inches(0.55),
             "招谁", size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.7), Inches(2.0), Inches(5.6), Inches(4.3),
             "优先：AI / 集成电路 / 智能机器人研发\n"
             "可组合：生物医药研发、汽车配套研发\n"
             "\n"
             "不招：整车 4S、整车展厅、汽车卖场\n"
             "先不承诺：公寓（C6 要先做合规）\n"
             "\n"
             "1F+2F 高空间做展示和接待，甲方自留或冠名",
             size=16, color=INK)

    add_rect(s, Inches(6.8), Inches(1.25), Inches(6.1), Inches(5.35), fill=CLOUD)
    add_rect(s, Inches(6.8), Inches(1.25), Inches(6.1), Inches(0.55), fill=GOLD)
    add_text(s, Inches(6.8), Inches(1.25), Inches(6.1), Inches(0.55),
             "怎么招", size=18, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(7.05), Inches(2.0), Inches(5.6), Inches(4.3),
             "1  建库，决策人到人\n"
             "2  拜访、带看，48 小时出纪要\n"
             "3  绿区直接报价，超线请示您\n"
             "4  统管中介，只留一套客户名单\n"
             "5  盯到盖章、保证金到账\n"
             "\n"
             "每周书面漏斗。出 PPT 不算完成。",
             size=16, color=INK)

    # 4 收费
    s = new_slide(prs)
    add_chrome(s, prs, page_no=4, label="收费",
               title="三笔账 · 招到才抽佣",
               subtitle="建议口径，商务简版确认栏可改数字。未税，增值税另计。")

    fees = [
        ("启动费", "15 万", "一次性 · 进场\n全额抵扣成功佣金"),
        ("月度费", "8 万 / 月", "干活 · 3–4 人组\n不含中介和活动硬成本"),
        ("成功佣金", "8% / 链主 12%", "首年租金为基数\n扣免租 · 没签成不抽"),
    ]
    for i, (t, n, b) in enumerate(fees):
        x = Inches(0.45) + Inches(4.2) * i
        add_rect(s, x, Inches(1.25), Inches(4.0), Inches(3.15), fill=CLOUD)
        add_rect(s, x, Inches(1.25), Inches(4.0), Inches(0.5),
                 fill=GOLD if i == 2 else NAVY)
        add_text(s, x, Inches(1.25), Inches(4.0), Inches(0.5), t,
                 size=16, bold=True,
                 color=NAVY if i == 2 else WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, Inches(1.85), Inches(4.0), Inches(0.85), n,
                 size=26, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.15), Inches(2.75), Inches(3.7), Inches(1.4),
                 b, size=15, color=INK, align=PP_ALIGN.CENTER)

    add_rect(s, Inches(0.45), Inches(4.6), Inches(12.4), Inches(2.0), fill=NAVY)
    add_text(s, Inches(0.7), Inches(4.75), Inches(12.0), Inches(1.7),
             "和空置比：4,000 ㎡ 空一年，6.5 元/㎡·天，大约少收 949 万。\n"
             "付给操盘方：月度 96 万 + 成功佣金约 57 万 ≈ 153 万。\n"
             "甲方自有未报备客户成交，不付成功佣金。中介佣金另付。",
             size=16, color=WHITE)

    # 5 签法
    s = new_slide(prs)
    add_chrome(s, prs, page_no=5, label="签法",
               title="可以先干 90 天",
               subtitle="怕 12 个月锁死，就先看人干活。收费标准两套签法一样。")

    add_table(
        s, Inches(0.45), Inches(1.25), Inches(12.4), Inches(3.4),
        ["", "方案 A · 12 个月独家（推荐）", "方案 B · 先 90 天"],
        [
            ["适合", "决心把楼租出去", "想先看执行再锁独家"],
            ["怎么停", "满 6 个月没有锚定意向，可停月度费", "90 天到期可以不转正"],
            ["合同章", "始终在您手里", "始终在您手里"],
            ["报价", "绿区我们直接报，黄区 48 小时您批", "同样"],
        ],
        header_size=14, body_size=14,
        col_widths=[Inches(1.8), Inches(5.3), Inches(5.3)],
    )

    add_rect(s, Inches(0.45), Inches(4.9), Inches(12.4), Inches(1.7), fill=CLOUD)
    add_text(s, Inches(0.7), Inches(5.05), Inches(12.0), Inches(1.4),
             "绿区：租金 ≥ 6.5 元/㎡·天 · 免租 ≤ 9 个月 · 装补 ≤ 600 元/㎡\n"
             "不承诺政府政策、落户、公寓。不收租户或中介回扣。\n"
             "详细条款在 05 号全稿，会上先把简版确认栏勾完。",
             size=16, color=INK)

    # 6 请定
    s = new_slide(prs)
    add_chrome(s, prs, page_no=6, label="拍板",
               title="今天勾这四项就够",
               subtitle="数字写在商务简版确认栏。法务随后出全稿，不改收费结构。")

    asks = [
        ("01", "签法", "A 十二个月独家，或 B 先九十天。"),
        ("02", "三个数", "启动费 / 月度费 / 成功佣金。可按 15 万、8 万、8%·12% 先填。"),
        ("03", "对接人", "唯一商务对接人 + 红区谁拍板。"),
        ("04", "绿区", "能否直接报价。每单都开会，节奏会断。"),
    ]
    for i, (n, k, v) in enumerate(asks):
        y = Inches(1.25) + Inches(1.25) * i
        add_round(s, Inches(0.5), y, Inches(1.15), Inches(1.05), n,
                  fill=GOLD, color=NAVY, size=20)
        add_text(s, Inches(1.9), y, Inches(2.2), Inches(1.05), k,
                 size=22, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(4.2), y, Inches(8.5), Inches(1.05), v,
                 size=18, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    total = len(SLIDES)
    for sl in SLIDES:
        for shape in sl.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text and run.text.strip().endswith(" / 0"):
                        run.text = run.text.replace(" / 0", f" / {total}")

    out = Path(__file__).resolve().parent.parent / "docs" / "advisory" / \
        "deck" / "冠松01楼-招商合作-会上版.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"✓ Deck written: {out}  ({total} slides)")
    return out


if __name__ == "__main__":
    build()
