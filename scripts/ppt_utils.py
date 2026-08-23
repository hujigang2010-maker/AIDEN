# -*- coding: utf-8 -*-
"""16:9 深蓝演示稿绘制函数。"""
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

NAVY = RGBColor(0x0B, 0x16, 0x2C)
NAVY2 = RGBColor(0x12, 0x24, 0x45)
CARD = RGBColor(0x16, 0x2C, 0x52)
BLUE = RGBColor(0x33, 0x70, 0xFF)
CYAN = RGBColor(0x2C, 0xD3, 0xC8)
GOLD = RGBColor(0xF0, 0xC0, 0x5A)
WHITE = RGBColor(0xF7, 0xFA, 0xFF)
MUTED = RGBColor(0xA8, 0xB8, 0xD0)
PINK = RGBColor(0xFF, 0x7A, 0x9A)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT = "微软雅黑"


def set_run_font(run, name=FONT, size=18, bold=False, color=WHITE):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {"typeface": name})
            rPr.append(el)
        else:
            el.set("typeface", name)
    latin = rPr.find(qn("a:latin"))
    if latin is None:
        latin = rPr.makeelement(qn("a:latin"), {"typeface": name})
        rPr.append(latin)
    else:
        latin.set("typeface", name)


def add_rect(slide, l, t, w, h, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def add_round(slide, l, t, w, h, fill, radius=0.08):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    return sh


def add_text(
    slide,
    l,
    t,
    w,
    h,
    text,
    size=18,
    bold=False,
    color=WHITE,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    if isinstance(text, list):
        lines = text
    else:
        lines = str(text).split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = line
        set_run_font(run, size=size, bold=bold, color=color)
    return tb


def add_paras(slide, l, t, w, h, items, size=16, color=WHITE, bullet=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_top = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        p.line_spacing = 1.15
        run = p.add_run()
        run.text = f"• {item}" if bullet else item
        set_run_font(run, size=size, color=color)
    return tb


def paint_bg(slide):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    add_rect(slide, Inches(0), Inches(0), Inches(0.14), SLIDE_H, BLUE)
    add_rect(slide, Inches(0), Inches(7.38), SLIDE_W, Inches(0.12), CYAN)


def footer(slide, page, total=14):
    add_text(
        slide,
        Inches(0.5),
        Inches(7.08),
        Inches(9.5),
        Inches(0.28),
        "飞书 AI Builder Demo Day #4  ·  观摩与转发材料  ·  非官方议程",
        size=11,
        color=MUTED,
    )
    add_text(
        slide,
        Inches(11.4),
        Inches(7.08),
        Inches(1.5),
        Inches(0.28),
        f"{page} / {total}",
        size=11,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def header(slide, kicker, title):
    add_text(slide, Inches(0.5), Inches(0.28), Inches(12.2), Inches(0.32), kicker, size=13, color=CYAN, bold=True)
    add_text(slide, Inches(0.5), Inches(0.58), Inches(12.2), Inches(0.5), title, size=28, color=WHITE, bold=True)
