#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成《AI × 不动产资产管理高级研修班》对外展示方案 PPT。

定位：面向合作方、渠道与意向学员的 16:9 汇报稿。
主干：AI+不动产培训；澳洲考察仅作圈层延展，不作为主线。

运行：python3 scripts/build_ai_real_estate_deck.py
输出：output/AI+不动产资产管理高级研修班_合作方案.pptx
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# 设计规范：不动产稳健蓝 + AI 青绿 + 圈层金
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
NAVY2 = RGBColor(0x14, 0x32, 0x56)
NAVY3 = RGBColor(0x1C, 0x42, 0x6E)
TEAL = RGBColor(0x0E, 0x8A, 0x7D)
TEAL_DEEP = RGBColor(0x0A, 0x6B, 0x61)
FUDAN = RGBColor(0x9E, 0x1B, 0x32)  # 复旦红：品牌与意识形态相关强调
GOLD = RGBColor(0xC4, 0xA3, 0x5A)
GOLD_PALE = RGBColor(0xF4, 0xEB, 0xD3)
CREAM = RGBColor(0xF6, 0xF3, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1E, 0x2A, 0x3A)
GRAY = RGBColor(0x5C, 0x6B, 0x7A)
MUTED = RGBColor(0x8A, 0x96, 0xA3)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xE4, 0xDC, 0xCC)
SOFT = RGBColor(0xEE, 0xF6, 0xF4)

FONT = "WenQuanYi Micro Hei"
FOOTER = "三方联合主办 · AI × 不动产资产管理高级研修班 · 2026 秋季首期"
HOSTS = [
    "复旦大学住房政策研究中心",
    "上海市杨浦区科技企业联合会",
    "上海市工商联房地产商会资产管理分会",
]

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT_DIR = os.path.join(ROOT, "output")
OUT = os.path.join(OUT_DIR, "AI+不动产资产管理高级研修班_合作方案.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def _set_font(run, size, color, bold=False, italic=False, name=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def rect(s, x, y, w, h, fill, line=None, line_w=None, round_=False, radius=0.06):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE, x, y, w, h
    )
    if round_:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w or 1)
    shp.shadow.inherit = False
    return shp


def txt(s, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True, spacing=1.08):
    """lines: tuple 或 list of tuples (text, size, color, bold[, space_after])。"""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = spacing
        segs = ln if isinstance(ln, list) else [ln]
        sa = None
        for seg in segs:
            text, size, color = seg[0], seg[1], seg[2]
            bold = seg[3] if len(seg) > 3 else False
            if len(seg) > 4 and seg[4] is not None:
                sa = seg[4]
            r = p.add_run()
            r.text = text
            _set_font(r, size, color, bold)
        if sa is not None:
            p.space_after = Pt(sa)
        p.space_before = Pt(0)
    return tb


def header(s, idx, kicker, title, subtitle=None):
    rect(s, 0, 0, SW, SH, CREAM)
    rect(s, 0, 0, Inches(0.14), SH, TEAL)
    rect(s, 0, 0, SW, Inches(0.08), NAVY)
    txt(
        s,
        Inches(0.55),
        Inches(0.28),
        Inches(10.4),
        Inches(0.28),
        [(FOOTER, 10, MUTED, False)],
    )
    txt(
        s,
        Inches(0.55),
        Inches(0.56),
        Inches(11.6),
        Inches(0.62),
        [[(kicker + "  ", 13, TEAL, True), (title, 24, NAVY, True)]],
    )
    if subtitle:
        txt(
            s,
            Inches(0.55),
            Inches(1.18),
            Inches(11.8),
            Inches(0.34),
            [(subtitle, 12.5, GRAY, False)],
        )
        bar_y = Inches(1.56)
    else:
        bar_y = Inches(1.28)
    rect(s, Inches(0.55), bar_y, Inches(12.2), Pt(2.2), GOLD)
    txt(
        s,
        Inches(12.15),
        Inches(7.08),
        Inches(0.85),
        Inches(0.28),
        [(f"{idx:02d}", 12, MUTED, True)],
        align=PP_ALIGN.RIGHT,
    )


def card(s, x, y, w, h, accent=TEAL):
    rect(s, x, y, w, h, CARD, line=LINE, line_w=1, round_=True, radius=0.05)
    rect(s, x, y, Inches(0.10), h, accent)


# ===========================================================================
# 01 封面
# ===========================================================================
def cover():
    s = slide()
    rect(s, 0, 0, SW, SH, NAVY)
    rect(s, 0, 0, SW, Inches(0.12), TEAL)
    rect(s, Inches(9.85), 0, Inches(3.48), SH, NAVY2)
    rect(s, Inches(10.35), Inches(1.35), Inches(2.45), Inches(2.45), TEAL)
    rect(s, Inches(10.85), Inches(1.85), Inches(1.45), Inches(1.45), GOLD)

    txt(s, Inches(0.85), Inches(1.05), Inches(8.8), Inches(0.4),
        [("合 作 方 案  ·  2026 秋 季 首 期", 14, GOLD, True)])
    txt(s, Inches(0.82), Inches(1.7), Inches(9.2), Inches(2.15),
        [[("AI × 不动产", 44, WHITE, True, 6)],
         [("资产管理高级研修班", 36, WHITE, True)]])
    txt(s, Inches(0.85), Inches(4.05), Inches(8.8), Inches(0.55),
        [("不动产做主干，AI 做素养，圈层做长期价值。", 17, RGBColor(0xC5, 0xD4, 0xE4), False)])
    txt(s, Inches(0.85), Inches(4.58), Inches(8.8), Inches(0.38),
        [("面向上海及长三角不动产从业者的高性价比研修与资源连接方案", 13, MUTED, False)])

    chips = [
        ("周期", "6 个月 · 每月 1 天"),
        ("班型", "首期 25–60 人"),
        ("定价", "2.0–2.48 万元"),
        ("开班", "2026 年 11 月"),
    ]
    x = Inches(0.85)
    for k, v in chips:
        rect(s, x, Inches(5.55), Inches(2.05), Inches(1.12), NAVY2, round_=True, radius=0.08)
        txt(s, x + Inches(0.16), Inches(5.68), Inches(1.75), Inches(0.85),
            [[(k, 11, GOLD, True)], [(v, 13, WHITE, True)]])
        x += Inches(2.2)

    txt(s, Inches(10.15), Inches(4.12), Inches(2.95), Inches(0.32),
        [("联 合 主 办", 12, GOLD, True)], align=PP_ALIGN.CENTER)
    hosts_short = [
        "复旦大学住房政策研究中心",
        "杨浦区科技企业联合会",
        "工商联房地产商会资产管理分会",
    ]
    yy = Inches(4.50)
    for i, h in enumerate(hosts_short):
        txt(s, Inches(10.05), yy, Inches(3.15), Inches(0.7),
            [[(f"{i+1}  ", 12, GOLD, True), (h, 11, WHITE, True)]],
            align=PP_ALIGN.LEFT)
        yy += Inches(0.72)


# ===========================================================================
# 02 目录
# ===========================================================================
def toc():
    s = slide()
    header(s, 2, "目录", "今天对齐什么", "一条主线讲清楚：为什么做、三家怎么分工、复旦品牌用到哪一层、钱怎么分。")
    items = [
        ("01", "机会判断", "地产进入资产管理时代，AI 正在改写岗位与工具。"),
        ("02", "项目定位", "不是纯 AI 课，而是不动产实操 + AI 素养 + 行业圈层。"),
        ("03", "产品设计", "六个月班型、课程地图、师资与社群机制。"),
        ("04", "主办与分成", "三家联合主办；复旦品牌分层决定意识形态审核与财务分成。"),
        ("05", "推进节奏", "9 月启动招生，11 月开班；先把培训做成入口。"),
    ]
    y = Inches(1.85)
    for no, title, desc in items:
        rect(s, Inches(0.55), y, Inches(12.2), Inches(0.92), CARD, line=LINE, line_w=1, round_=True)
        rect(s, Inches(0.55), y, Inches(1.15), Inches(0.92), NAVY, round_=True)
        txt(s, Inches(0.55), y, Inches(1.15), Inches(0.92),
            [(no, 18, GOLD, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, Inches(1.95), y + Inches(0.12), Inches(10.4), Inches(0.7),
            [[(title, 18, NAVY, True, 2)], [(desc, 13, GRAY, False)]])
        y += Inches(1.02)


# ===========================================================================
# 03 机会判断
# ===========================================================================
def opportunity():
    s = slide()
    header(s, 3, "01  机会判断", "为什么是现在，为什么是不动产",
           "行业在换赛道，人在换岗位。培训要接住这两股力，而不是再讲一遍通用 AI。")

    cols = [
        (NAVY, "行业侧", "从开发销售转向资产管理",
         ["存量时代，持有、运营、处置成为主业。",
          "不良资产、REITs、资本化退出成为刚需话题。",
          "纯开发逻辑的课已经很难卖，资产管理才有人买单。"]),
        (TEAL, "技术侧", "AI 改写工具，但不改行业本质",
         ["中年从业者有转型压力，但听不懂纯技术课。",
          "市场上 AI 课同质化严重，听半小时就触顶。",
          "真正能成交的，是「AI 用在不动产上能干什么」。"]),
        (RGBColor(0x8A, 0x6A, 0x2A), "需求侧", "买的是圈层，不是课件",
         ["EMBA 动辄几十万，周期长、门槛高。",
          "两万元左右换一个可走动的行业圈子，更贴近当下。",
          "上海及长三角商会会员、国企、银行是现成客源。"]),
    ]
    x = Inches(0.5)
    for accent, kicker, title, bullets in cols:
        rect(s, x, Inches(1.85), Inches(3.95), Inches(4.85), CARD, line=LINE, line_w=1, round_=True)
        rect(s, x, Inches(1.85), Inches(3.95), Inches(1.35), accent, round_=True)
        rect(s, x, Inches(2.75), Inches(3.95), Inches(0.45), accent)
        txt(s, x + Inches(0.22), Inches(1.98), Inches(3.5), Inches(1.1),
            [[(kicker, 12, GOLD_PALE, True, 2)], [(title, 16, WHITE, True)]])
        yy = Inches(3.4)
        for b in bullets:
            txt(s, x + Inches(0.28), yy, Inches(3.4), Inches(0.95),
                [[("●  ", 12, TEAL, True), (b, 13, INK, False)]])
            yy += Inches(1.0)
        x += Inches(4.15)


# ===========================================================================
# 04 项目定位
# ===========================================================================
def positioning():
    s = slide()
    header(s, 4, "02  项目定位", "一句话讲清楚这门课",
           "对外可以说 AI，对内必须守住不动产。否则既听不懂，也卖不掉。")

    rect(s, Inches(0.5), Inches(1.85), Inches(12.3), Inches(1.35), NAVY, round_=True)
    txt(s, Inches(0.75), Inches(2.0), Inches(11.8), Inches(1.05),
        [[("不动产资产管理是主干，AI 只做基础素养。", 20, WHITE, True, 4)],
         [("班集体、学员互访和长期圈层，才是复购与转介绍的真正产品。", 14, RGBColor(0xC5, 0xD4, 0xE4), False)]],
        anchor=MSO_ANCHOR.MIDDLE)

    left = [
        ("不是", "一场通用 AI 公开课", "不堆模型名词，不做成科技展讲解。"),
        ("不是", "几十万级商学院项目", "周期更短、价格更友好、行业更垂直。"),
        ("不是", "一次听完就散的沙龙", "六个月同班，才能把人变成可持续连接。"),
    ]
    right = [
        ("是", "不动产实操研修", "资产管理、不良资产、REITs，对接本行工作。"),
        ("是", "高性价比行业圈层", "用两万元左右进入可走动的不动产资源网。"),
        ("是", "后续业务的流量入口", "先把班做成品牌，再谈科创对接与 FA。"),
    ]
    y = Inches(3.4)
    for (a1, t1, d1), (a2, t2, d2) in zip(left, right):
        rect(s, Inches(0.5), y, Inches(6.0), Inches(1.12), CARD, line=LINE, line_w=1, round_=True)
        rect(s, Inches(0.5), y, Inches(0.9), Inches(1.12), RGBColor(0x8B, 0x3A, 0x3A), round_=True)
        txt(s, Inches(0.5), y, Inches(0.9), Inches(1.12),
            [(a1, 14, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, Inches(1.6), y + Inches(0.14), Inches(4.7), Inches(0.85),
            [[(t1, 15, NAVY, True, 2)], [(d1, 12, GRAY, False)]])

        rect(s, Inches(6.8), y, Inches(6.0), Inches(1.12), CARD, line=LINE, line_w=1, round_=True)
        rect(s, Inches(6.8), y, Inches(0.9), Inches(1.12), TEAL, round_=True)
        txt(s, Inches(6.8), y, Inches(0.9), Inches(1.12),
            [(a2, 14, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, Inches(7.9), y + Inches(0.14), Inches(4.7), Inches(0.85),
            [[(t2, 15, NAVY, True, 2)], [(d2, 12, GRAY, False)]])
        y += Inches(1.22)


# ===========================================================================
# 05 核心洞察
# ===========================================================================
def insight():
    s = slide()
    header(s, 5, "02  项目定位", "学员真正在买什么",
           "招生话术可以讲 AI，成交逻辑要讲圈层与转型。")

    rect(s, Inches(0.5), Inches(1.82), Inches(12.3), Inches(1.55), SOFT, round_=True)
    txt(s, Inches(0.8), Inches(1.98), Inches(11.7), Inches(1.25),
        [[("核心判断", 12, TEAL, True, 4)],
         [("多数人出来上课，不是来把 AI 学成专家，而是用可接受的成本进入一个高质量行业圈子，顺便补齐 AI 时代的常识。",
           16, NAVY, True)]],
        spacing=1.15)

    items = [
        ("低成本替代 EMBA", "长江商学院一类项目一上来就是几十万、周期长。本班用约 2 万元，换六个月可持续走动的同业关系。"),
        ("中年转型刚需", "付费主力是有支付能力、又面临岗位迭代的中年从业者，不是应届学生。"),
        ("行业垂直才成交", "客源来自房地产商会、国企、银行支行行长，以及广告、财税、律师等「百搭」上下游。"),
        ("AI 是入场券，不是主菜", "半小时讲清工具地图即可。再深，听不懂；只讲地产、不带 AI，又缺当下话题。"),
    ]
    positions = [
        (Inches(0.5), Inches(3.55)),
        (Inches(6.8), Inches(3.55)),
        (Inches(0.5), Inches(5.2)),
        (Inches(6.8), Inches(5.2)),
    ]
    for (title, body), (x, y) in zip(items, positions):
        rect(s, x, y, Inches(6.0), Inches(1.5), CARD, line=LINE, line_w=1, round_=True)
        rect(s, x, y, Inches(0.10), Inches(1.5), GOLD)
        txt(s, x + Inches(0.28), y + Inches(0.14), Inches(5.5), Inches(1.22),
            [[(title, 15, NAVY, True, 4)], [(body, 12.5, GRAY, False)]])


# ===========================================================================
# 06 目标客群
# ===========================================================================
def audience():
    s = slide()
    header(s, 6, "02  项目定位", "谁适合来，从哪里来",
           "先吃透上海及长三角不动产圈层，再谈扩科创、做出海。")

    personas = [
        ("核心盘", TEAL, "房地产商会会员", "上海本地会员 + 苏浙皖异地商会。每月来一次上海，可接受。"),
        ("支付盘", NAVY, "中年转型从业者", "有学费支付能力，对 AI 焦虑，需要新圈子和新话术。"),
        ("信用盘", NAVY3, "国企 / 银行支行行长", "不便入会，但需要资产、不良与客户资源对接。"),
        ("连接盘", RGBColor(0x8A, 0x6A, 0x2A), "地产上下游「百搭」", "广告、信息技术、财税、律师、二房东、社区商业。"),
    ]
    x = Inches(0.5)
    for tag, color, title, desc in personas:
        rect(s, x, Inches(1.85), Inches(3.0), Inches(3.15), CARD, line=LINE, line_w=1, round_=True)
        rect(s, x, Inches(1.85), Inches(3.0), Inches(0.7), color, round_=True)
        txt(s, x, Inches(1.85), Inches(3.0), Inches(0.7),
            [(tag, 13, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, x + Inches(0.18), Inches(2.7), Inches(2.64), Inches(2.1),
            [[(title, 16, NAVY, True, 8)], [(desc, 13, GRAY, False)]])
        x += Inches(3.2)

    rect(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.5), NAVY, round_=True)
    txt(s, Inches(0.8), Inches(5.38), Inches(11.7), Inches(1.15),
        [[("获客原则", 12, GOLD, True, 4)],
         [("自有会员给面子、异地商协会秘书长做渠道（返佣 20%–25%）、国企与银行点对点邀请。外地客源不必强拉入会，用班级把人留下即可。",
           14, WHITE, False)]],
        spacing=1.12)


# ===========================================================================
# 07 产品形态
# ===========================================================================
def product():
    s = slide()
    header(s, 7, "03  产品设计", "班型怎么开",
           "沿用已跑通的「六个月、每月一天、校内上课」标准配置，降低试错成本。")

    specs = [
        ("周期", "6 个月", "从 5 个月拉长，给社群和互访留足时间。"),
        ("课时", "每月 1 天", "周五或周末集中授课，便于外地学员往返。"),
        ("场地", "高校校内", "证书、发票、场地一次性对齐，观感更稳。"),
        ("规模", "25–60 人", "25 人起班，60 人更利于圈层生态；首期建议冲 40–60。"),
        ("滚动", "每两月一期", "不是一次性项目，做成可复制的班级产品。"),
        ("班费", "5–10 万元", "从学费中预留，用于互访、聚餐和班委活动。"),
    ]
    x0, y0 = Inches(0.5), Inches(1.85)
    cw, ch = Inches(4.0), Inches(2.35)
    for i, (k, v, d) in enumerate(specs):
        r, c = divmod(i, 3)
        x = x0 + c * Inches(4.15)
        y = y0 + r * Inches(2.5)
        rect(s, x, y, cw, ch, CARD, line=LINE, line_w=1, round_=True)
        txt(s, x + Inches(0.28), y + Inches(0.22), Inches(3.45), Inches(0.35),
            [(k, 12, TEAL, True)])
        txt(s, x + Inches(0.28), y + Inches(0.55), Inches(3.45), Inches(0.5),
            [(v, 22, NAVY, True)])
        rect(s, x + Inches(0.28), y + Inches(1.15), Inches(0.7), Pt(2.5), GOLD)
        txt(s, x + Inches(0.28), y + Inches(1.35), Inches(3.45), Inches(0.8),
            [(d, 13, GRAY, False)])


# ===========================================================================
# 08 课程地图
# ===========================================================================
def curriculum():
    s = slide()
    header(s, 8, "03  产品设计", "六个月课程地图",
           "每月一天主课 + 贯穿始终的学员单位互访。AI 只占必要篇幅，不动产实操占主干。")

    modules = [
        ("01", "开班与 AI 素养", "工具地图、场景入门、班委组建、破冰互访。"),
        ("02", "不动产资产管理", "资产盘点、持有运营、现金流与估值框架。"),
        ("03", "不良资产处置", "识别、尽调、重组、司法与市场化路径。"),
        ("04", "REITs 与资本化", "公募 REITs、证券化、退出与结构安排。"),
        ("05", "AI 赋能资产运营", "租赁、招商、客服、研报的轻量落地，不深挖算法。"),
        ("06", "结业路演对接", "小组课题、资源对接、结业仪式、同学会启动。"),
    ]
    x0, y0 = Inches(0.5), Inches(1.82)
    cw, ch = Inches(4.0), Inches(2.32)
    for i, (no, title, desc) in enumerate(modules):
        r, c = divmod(i, 3)
        x = x0 + c * Inches(4.15)
        y = y0 + r * Inches(2.48)
        rect(s, x, y, cw, ch, CARD, line=LINE, line_w=1, round_=True)
        rect(s, x, y, Inches(0.72), ch, NAVY if i < 4 else TEAL)
        txt(s, x, y, Inches(0.72), ch,
            [(no, 16, GOLD, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, x + Inches(0.9), y + Inches(0.35), Inches(2.9), Inches(1.7),
            [[(title, 16, NAVY, True, 8)], [(desc, 13, GRAY, False)]])


# ===========================================================================
# 09 师资与资源
# ===========================================================================
def faculty():
    s = slide()
    header(s, 9, "03  产品设计", "谁来讲，凭什么相信",
           "师资成本可控：以头部企业高管分享 + 不动产实操老师为主，不堆高价商学院牌。")

    left_items = [
        ("AI 侧 · 科企联入口", "腾讯、阿里，以及宇树、智元机器人等参展企业高管。杨浦科企联对接科创师资，讲产业趋势与可用工具。"),
        ("不动产侧 · 商会入口", "资产管理、不良处置、REITs 实操老师，由房商会资管分会组织本行师资。这是续班和转介绍的根基。"),
        ("学术侧 · 复旦中心", "课纲审定、校内场地、结业路径。用中心名义还是学校证书，决定审核深度和分成深度。"),
    ]
    y = Inches(1.85)
    for title, body in left_items:
        rect(s, Inches(0.5), y, Inches(7.55), Inches(1.5), CARD, line=LINE, line_w=1, round_=True)
        rect(s, Inches(0.5), y, Inches(0.12), Inches(1.5), TEAL)
        txt(s, Inches(0.85), y + Inches(0.22), Inches(6.95), Inches(1.1),
            [[(title, 16, NAVY, True, 4)], [(body, 13, GRAY, False)]])
        y += Inches(1.62)

    rect(s, Inches(8.25), Inches(1.85), Inches(4.55), Inches(4.85), NAVY, round_=True)
    txt(s, Inches(8.5), Inches(2.1), Inches(4.05), Inches(0.4),
        [("课堂上要守住的边界", 14, GOLD, True)])
    notes = [
        "AI 讲皮毛即可，半小时能讲清工具地图。",
        "太深听不懂，全是 AI 又听腻。",
        "不动产课能不能卖，取决于客群是否垂直。",
        "挂复旦名义的文案、讲师、涉外内容，须先过中心审。",
        "讲师费用按「认圈子、滚动流量」来谈，不按商学院时薪来堆。",
    ]
    yy = Inches(2.65)
    for n in notes:
        txt(s, Inches(8.55), yy, Inches(4.0), Inches(0.7),
            [[("▸  ", 13, TEAL, True), (n, 13, WHITE, False)]])
        yy += Inches(0.72)


# ===========================================================================
# 10 社群运营
# ===========================================================================
def community():
    s = slide()
    header(s, 10, "03  产品设计", "班集体怎么运转",
           "课只是由头。真正把人留住的，是班委、互访和滚动开班形成的圈层密度。")

    steps = [
        ("1", "开班建制", "设班长、学习委员、组长。把 40–60 人拆成可走动的小组。"),
        ("2", "每月主课", "校内一天集中授课，统一节奏，便于外地学员安排。"),
        ("3", "学员互访", "课间穿插单位参访，把课堂关系变成业务关系。"),
        ("4", "班费活动", "预留 5–10 万元做聚餐、走访和小型沙龙，不另收费。"),
        ("5", "滚动续班", "每两个月新开一期，老学员可回流当链接人。"),
        ("6", "出口对接", "结业后筛种子选手，进入科创对接与考察名单。"),
    ]
    x0, y0 = Inches(0.5), Inches(1.85)
    for i, (no, title, desc) in enumerate(steps):
        r, c = divmod(i, 3)
        x = x0 + c * Inches(4.15)
        y = y0 + r * Inches(2.45)
        rect(s, x, y, Inches(4.0), Inches(2.25), CARD, line=LINE, line_w=1, round_=True)
        oval = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.25), y + Inches(0.28), Inches(0.48), Inches(0.48))
        oval.fill.solid()
        oval.fill.fore_color.rgb = TEAL if i < 4 else GOLD
        oval.line.fill.background()
        oval.shadow.inherit = False
        tf = oval.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r0 = p.add_run()
        r0.text = no
        _set_font(r0, 14, WHITE, True)
        txt(s, x + Inches(0.88), y + Inches(0.32), Inches(2.9), Inches(0.42),
            [(title, 16, NAVY, True)], anchor=MSO_ANCHOR.MIDDLE)
        txt(s, x + Inches(0.25), y + Inches(0.95), Inches(3.5), Inches(1.05),
            [(desc, 13, GRAY, False)])


# ===========================================================================
# 11 三家主办
# ===========================================================================
def hosts():
    s = slide()
    header(s, 11, "04  主办与分成", "三家联合主办，各管一段",
           "牌子可以并列，权责必须拆开。复旦管名义与审核，科企联管科创入口，商会管不动产圈层。")

    cols = [
        (FUDAN, "复旦大学住房政策研究中心", "学术审核 · 品牌边界",
         ["课纲、讲师、招生文案审定", "校内场地与结业路径", "中心非独立法人，学校是主管单位",
          "名义用得越深，审核越严、分成越高"]),
        (TEAL, "上海市杨浦区科技企业联合会", "科创入口 · B 端资源",
         ["对接 AI / 机器人等企业师资", "杨浦园区、企业参访与赞助", "后续科创大赛、特展的主场",
          "适合把班级活动打包给产业园"]),
        (NAVY, "上海市工商联房地产商会资产管理分会", "客源盘 · 圈层运营",
         ["上海及长三角不动产会员", "国企、银行与上下游「百搭」", "学员互访、班委、班费活动",
          "核心招生与转介绍出口"]),
    ]
    x = Inches(0.45)
    for color, name, role, pts in cols:
        rect(s, x, Inches(1.82), Inches(4.05), Inches(4.45), CARD, line=LINE, line_w=1, round_=True)
        rect(s, x, Inches(1.82), Inches(4.05), Inches(1.28), color, round_=True)
        rect(s, x, Inches(2.72), Inches(4.05), Inches(0.38), color)
        txt(s, x + Inches(0.12), Inches(1.90), Inches(3.8), Inches(1.1),
            [[(name, 13, WHITE, True, 2)], [(role, 12, GOLD_PALE, True)]],
            align=PP_ALIGN.CENTER)
        yy = Inches(3.28)
        for p in pts:
            txt(s, x + Inches(0.2), yy, Inches(3.65), Inches(0.68),
                [[("●  ", 12, GOLD, True), (p, 13, INK, False)]])
            yy += Inches(0.7)
        x += Inches(4.2)

    rect(s, Inches(0.45), Inches(6.42), Inches(12.4), Inches(0.58), SOFT, round_=True)
    txt(s, Inches(0.65), Inches(6.48), Inches(12.0), Inches(0.46),
        [("原则：三家都上主办栏，但复旦的「品牌 / 意识形态」和「学校财务分成」是两笔账，必须分开谈，不能混成一句「给复旦一点」。",
          13, NAVY, True)],
        anchor=MSO_ANCHOR.MIDDLE)


# ===========================================================================
# 12 复旦品牌与意识形态
# ===========================================================================
def fudan_brand():
    s = slide()
    header(s, 12, "04  主办与分成", "复旦品牌不是免费的：名义、意识形态、分成绑在一起",
           "用中心名义是一层成本，用学校证书是另一层。前者可谈项目费，后者是学校制度分成。")

    levels = [
        ("L1 建议默认", TEAL, "中心联合主办",
         "对外写「住房政策研究中心联合主办」。三方联合结业，不开发票到学校，不用校徽做主视觉。"),
        ("L2 学员点选", FUDAN, "学校证书 + 学校发票",
         "学费进学校账户，发复旦体系官方结业证。审核、收费上限、安全责任全部按学校非学历培训走。"),
        ("L0 备选", GRAY, "不用复旦名义",
         "仅科企联 + 资管分会上主办。品牌最弱，也最省审核。招生话术会明显变软，一般不作为主推。"),
    ]
    y = Inches(1.78)
    for tag, color, title, body in levels:
        rect(s, Inches(0.45), y, Inches(7.55), Inches(1.18), CARD, line=LINE, line_w=1, round_=True)
        rect(s, Inches(0.45), y, Inches(1.55), Inches(1.18), color, round_=True)
        txt(s, Inches(0.45), y, Inches(1.55), Inches(1.18),
            [(tag, 12, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, Inches(2.15), y + Inches(0.12), Inches(5.65), Inches(0.95),
            [[(title, 15, NAVY, True, 2)], [(body, 12, GRAY, False)]])
        y += Inches(1.26)

    rect(s, Inches(8.2), Inches(1.78), Inches(4.65), Inches(4.95), FUDAN, round_=True)
    txt(s, Inches(8.42), Inches(1.95), Inches(4.25), Inches(0.4),
        [("意识形态与品牌红线", 14, GOLD, True)])
    rules = [
        "不得使用校名、校徽做招生主视觉。",
        "不得暗示学历、学位或复旦学分。",
        "固定表述：住房政策研究中心联合主办。",
        "课纲、讲师、海报、涉外内容须中心审定。",
        "投资、移民、收益承诺不得挂复旦名义。",
        "中心会过问安全：不合规住宿、无保险不办。",
        "学校账户收费有金额上限，机票酒店宜外置。",
    ]
    yy = Inches(2.45)
    for r in rules:
        txt(s, Inches(8.45), yy, Inches(4.2), Inches(0.55),
            [[("▸  ", 12, GOLD, True), (r, 12.5, WHITE, False)]])
        yy += Inches(0.58)


# ===========================================================================
# 13 两种分成结构
# ===========================================================================
def models():
    s = slide()
    header(s, 13, "04  主办与分成", "品牌用到哪一层，钱就按哪一层分",
           "L1 给中心品牌使用费；L2 给学校制度分成。两笔不能当成同一件事来谈。")

    cards = [
        (TEAL, "模式 A  ·  L1 中心名义（建议默认）",
         ["学费进第三方账户，不占用学校抬头。",
          "复旦中心：品牌使用 + 课纲审定 + 场地协调，建议学费的 10%–15%（待中心书面确认）。",
          "渠道返佣 20%–25%、班费 5–10 万先扣。",
          "剩余可分配收益：科企联 30% / 资管分会 35% / 运营执行 35%。",
          "结业证由三方联合颁发，不写学校学历口径。"]),
        (FUDAN, "模式 B  ·  L2 学校证书（仅点选）",
         ["学费进学校账户，发官方结业证、开学校发票。",
          "复旦学校：非学历培训制度分成，待书面确认；按同类高校 50%–80% 预留（交大口径约 80%）。",
          "学校过问内容、安全、收费上限；意识形态审核最严。",
          "剩余再按科企联 30% / 资管分会 35% / 运营 35% 切。",
          "仅当学员明确要「学校证书 / 学校发票」时启用。"]),
    ]
    xs = [Inches(0.45), Inches(6.85)]
    for (clr, title, pts), x in zip(cards, xs):
        rect(s, x, Inches(1.82), Inches(6.0), Inches(5.0), CARD, line=LINE, line_w=1, round_=True)
        rect(s, x, Inches(1.82), Inches(6.0), Inches(0.9), clr, round_=True)
        rect(s, x, Inches(2.32), Inches(6.0), Inches(0.4), clr)
        txt(s, x + Inches(0.22), Inches(1.92), Inches(5.55), Inches(0.7),
            [(title, 14, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
        yy = Inches(2.88)
        for p in pts:
            txt(s, x + Inches(0.25), yy, Inches(5.5), Inches(0.72),
                [[("●  ", 12, GOLD, True), (p, 12.5, INK, False)]])
            yy += Inches(0.74)


# ===========================================================================
# 12 收费与渠道
# ===========================================================================
def pricing():
    s = slide()
    header(s, 14, "04  主办与分成", "怎么报价，怎么分账",
           "对外给一个干净的价格带；对内把渠道、班费、复旦品牌费、三方可分配一次说清。")

    # 价格带
    bands = [
        ("入门沟通价", "约 2.0 万元", "用于试探与关系户面子价，不作为对外主报价。"),
        ("标准对外价", "2.48 万元", "主推档。覆盖讲师、场地、班费与渠道空间。"),
        ("渠道结算", "返佣 20%–25%", "异地商协会秘书长、合作机构按此区间谈。"),
    ]
    x = Inches(0.5)
    accents = [NAVY, TEAL, RGBColor(0x8A, 0x6A, 0x2A)]
    for (k, v, d), acc in zip(bands, accents):
        rect(s, x, Inches(1.85), Inches(4.0), Inches(2.15), CARD, line=LINE, line_w=1, round_=True)
        txt(s, x + Inches(0.25), Inches(2.02), Inches(3.5), Inches(0.32),
            [(k, 12, acc, True)])
        txt(s, x + Inches(0.25), Inches(2.35), Inches(3.5), Inches(0.5),
            [(v, 22, NAVY, True)])
        txt(s, x + Inches(0.25), Inches(2.95), Inches(3.5), Inches(0.8),
            [(d, 13, GRAY, False)])
        x += Inches(4.15)

    notes = [
        ("班级活动经费", "学费中预留 5–10 万元，专款用于互访与班级活动，不另向学员加收。"),
        ("报价原则", "当前高端商学院招生难，本班主打高性价比圈层，不与几十万项目正面比课时。"),
        ("不对外承诺", "不承诺投资收益、移民结果或高校学历；证书只写三方联合结业，不写复旦学历口径。"),
        ("渠道纪律", "返佣只走约定渠道。复旦品牌费不从渠道返佣里挤，单独列项。"),
    ]
    y = Inches(4.2)
    for i, (k, v) in enumerate(notes):
        c = i % 2
        r = i // 2
        x = Inches(0.5) + c * Inches(6.35)
        yy = y + r * Inches(1.2)
        rect(s, x, yy, Inches(6.15), Inches(1.08), CARD, line=LINE, line_w=1, round_=True)
        txt(s, x + Inches(0.22), yy + Inches(0.12), Inches(5.7), Inches(0.85),
            [[(k, 14, NAVY, True, 2)], [(v, 12.5, GRAY, False)]])


# ===========================================================================
# 13 收益测算
# ===========================================================================
def economics():
    s = slide()
    header(s, 15, "04  主办与分成", "模式 A 三方分成测算（示意）",
           "按 L1 中心名义、复旦品牌费 12%、渠道 22.5% 中位测算。数字用于三家对齐，待中心书面确认。")

    rows = [
        ("项目", "轻量班 25 人 × 2.00 万", "生态班 60 人 × 2.48 万"),
        ("学费收入", "50.0 万元", "148.8 万元"),
        ("渠道返佣 22.5%", "11.3 万元", "33.5 万元"),
        ("班级活动经费", "5.0 万元", "10.0 万元"),
        ("复旦中心品牌使用费 12%", "6.0 万元", "17.9 万元"),
        ("可分配收益", "27.7 万元", "87.4 万元"),
        ("杨浦科企联 30%", "8.3 万元", "26.2 万元"),
        ("资管分会 35%", "9.7 万元", "30.6 万元"),
        ("运营执行 35%", "9.7 万元", "30.6 万元"),
    ]
    col_w = [Inches(4.5), Inches(3.9), Inches(3.9)]
    row_h = Inches(0.46)
    x0, y0 = Inches(0.5), Inches(1.78)
    for r, row in enumerate(rows):
        y = y0 + r * row_h
        x = x0
        for c, cell in enumerate(row):
            w = col_w[c]
            if r == 0:
                fill = NAVY if c == 0 else TEAL
                color, bold = WHITE, True
            elif r == 4:
                fill = RGBColor(0xF8, 0xEB, 0xEE)
                color, bold = FUDAN, True
            elif r == 5:
                fill = SOFT
                color, bold = TEAL_DEEP, True
            elif r in (6, 7, 8):
                fill = CARD
                color, bold = INK, True
            else:
                fill = CARD
                color, bold = INK, False
            rect(s, x, y, w, row_h, fill, line=LINE, line_w=1)
            txt(s, x + Inches(0.12), y, w - Inches(0.2), row_h,
                [(cell, 12.5, color, bold)],
                anchor=MSO_ANCHOR.MIDDLE,
                align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)
            x += w

    rect(s, Inches(0.5), Inches(6.05), Inches(12.3), Inches(1.05), FUDAN, round_=True)
    txt(s, Inches(0.75), Inches(6.18), Inches(11.85), Inches(0.8),
        [[("模式 B 对照：", 13, GOLD, True),
          ("若学费进学校且按交大口径 80% 提取，60 人班学校拿走约 119 万，剩余约 30 万还要覆盖渠道与班费，三家几乎无账可分。故 L2 只作学员点选，默认走 L1。复旦学校确切比例以中心书面回复为准。",
           13, WHITE, False)]])


# ===========================================================================
# 14 三层生态
# ===========================================================================
def ecosystem():
    s = slide()
    header(s, 16, "05  推进节奏", "培训只是入口，后面还有两层",
           "战略可以画三层，执行必须先把第一层做成。至少跑满三期，再谈大赛与 FA。")

    layers = [
        ("第一层  ·  现在就做", TEAL, "AI+不动产培训",
         ["精准行业用户与班级品牌", "学费 + 渠道 + 班费可核算", "为后续两层提供名单与信任"]),
        ("第二层  ·  三期之后", NAVY, "长三角科创对接",
         ["从学员中筛 100–200 名种子", "特展约 4 万 / 场，系列约 50 万", "外地企业来沪参访 1–3 万 / 次"]),
        ("第三层  ·  偶得收益", GOLD, "FA 撮合佣金",
         ["只在有合作意向时抽取", "成功率低，不可当主营", "是补充，不是预算支柱"]),
    ]
    x = Inches(0.5)
    for tag, color, title, pts in layers:
        rect(s, x, Inches(1.85), Inches(4.0), Inches(3.95), CARD, line=LINE, line_w=1, round_=True)
        rect(s, x, Inches(1.85), Inches(4.0), Inches(1.2), color, round_=True)
        rect(s, x, Inches(2.6), Inches(4.0), Inches(0.45), color)
        txt(s, x + Inches(0.22), Inches(1.98), Inches(3.55), Inches(0.95),
            [[(tag, 12, GOLD_PALE if color != GOLD else NAVY, True, 2)],
             [(title, 18, WHITE if color != GOLD else NAVY, True)]])
        yy = Inches(3.25)
        for p in pts:
            txt(s, x + Inches(0.28), yy, Inches(3.45), Inches(0.7),
                [[("●  ", 12, TEAL, True), (p, 13.5, INK, False)]])
            yy += Inches(0.78)
        x += Inches(4.15)

    rect(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.7), SOFT, round_=True)
    txt(s, Inches(0.75), Inches(6.1), Inches(11.8), Inches(0.5),
        [("圈层延展（可选）：澳大利亚商务考察可作为高净值学员的深度连接产品，不并入本班主报价，单独成团、单独核算。",
          13, NAVY, False)],
        anchor=MSO_ANCHOR.MIDDLE)


# ===========================================================================
# 15 时间表
# ===========================================================================
def timeline():
    s = slide()
    header(s, 17, "05  推进节奏", "9–11 月黄金档，按周推进",
           "招生要两个月窗口。9 月中下旬启动，11 月开班；复旦品牌层级不定，文案就不能对外。")

    phases = [
        ("即刻", "方案对齐", ["确认 L1 中心名义为默认", "书面确认中心品牌费 10%–15%", "三方分成 30 / 35 / 35"]),
        ("9 月中下旬", "启动招生", ["固定表述：中心联合主办", "课纲讲师送中心审定", "商会点对点 + 科企联渠道"]),
        ("10 月", "满班冲刺", ["目标 40–60 人", "班委候选人预热", "首月课表与互访名单"]),
        ("11 月", "正式开班", ["三方联合开班仪式", "班费账户与互访节奏", "同步筹备第二期名单"]),
    ]
    x = Inches(0.5)
    colors = [TEAL, NAVY, NAVY3, GOLD]
    for (when, title, pts), color in zip(phases, colors):
        rect(s, x, Inches(1.85), Inches(3.0), Inches(4.85), CARD, line=LINE, line_w=1, round_=True)
        rect(s, x, Inches(1.85), Inches(3.0), Inches(1.35), color, round_=True)
        rect(s, x, Inches(2.75), Inches(3.0), Inches(0.45), color)
        txt(s, x + Inches(0.15), Inches(2.0), Inches(2.7), Inches(1.05),
            [[(when, 12, GOLD_PALE if color != GOLD else NAVY, True, 2)],
             [(title, 18, WHITE if color != GOLD else NAVY, True)]],
            align=PP_ALIGN.CENTER)
        yy = Inches(3.4)
        for p in pts:
            txt(s, x + Inches(0.2), yy, Inches(2.6), Inches(0.85),
                [[("●  ", 12, TEAL, True), (p, 13, INK, False)]])
            yy += Inches(0.95)
        x += Inches(3.2)


# ===========================================================================
# 16 下一步
# ===========================================================================
def next_steps():
    s = slide()
    header(s, 18, "05  推进节奏", "今天需要拍板的四件事",
           "方案可以继续改细节，但下面四项不定，招生文案就开不了口。")

    decisions = [
        ("01", "复旦层级", "默认 L1：住房政策研究中心联合主办，不用校徽、不开学校发票。"),
        ("02", "分成口径", "中心品牌费 10%–15%（待书面确认）；剩余科企联 30% / 分会 35% / 运营 35%。"),
        ("03", "定价班型", "对外 2.48 万元；首期按 40–60 人招，25 人为开班底线。"),
        ("04", "时间", "9 月中下旬启动招生，11 月开班，每两月滚动一期。"),
    ]
    y = Inches(1.82)
    for no, title, desc in decisions:
        rect(s, Inches(0.5), y, Inches(8.35), Inches(1.12), CARD, line=LINE, line_w=1, round_=True)
        rect(s, Inches(0.5), y, Inches(1.05), Inches(1.12), NAVY, round_=True)
        txt(s, Inches(0.5), y, Inches(1.05), Inches(1.12),
            [(no, 16, GOLD, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, Inches(1.75), y + Inches(0.16), Inches(6.85), Inches(0.82),
            [[(title, 16, NAVY, True, 2)], [(desc, 13, GRAY, False)]])
        y += Inches(1.2)

    rect(s, Inches(9.05), Inches(1.82), Inches(3.75), Inches(4.85), NAVY, round_=True)
    txt(s, Inches(9.28), Inches(2.05), Inches(3.3), Inches(0.45),
        [("会后立即分头做", 14, GOLD, True)])
    todos = [
        "请中心书面确认品牌费与用名",
        "送审课纲、讲师、招生主视觉",
        "列出首批 80 人邀请名单",
        "锁定第三方开票主体与场地",
        "准备一页招生海报与报名表",
    ]
    yy = Inches(2.6)
    for t in todos:
        txt(s, Inches(9.3), yy, Inches(3.25), Inches(0.7),
            [[("▸  ", 13, TEAL, True), (t, 13, WHITE, False)]])
        yy += Inches(0.72)


# ===========================================================================
# 17 结束页
# ===========================================================================
def closing():
    s = slide()
    rect(s, 0, 0, SW, SH, NAVY)
    rect(s, 0, 0, SW, Inches(0.12), TEAL)
    rect(s, 0, Inches(5.55), SW, Inches(0.06), GOLD)

    txt(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(0.4),
        [("先把一个班做成品牌，再谈后面的生态。", 16, GOLD, True)])
    txt(s, Inches(0.9), Inches(2.25), Inches(11.5), Inches(1.6),
        [[("不动产做主干", 32, WHITE, True, 8)],
         [("AI 做素养，圈层做长期价值", 28, WHITE, True)]])
    txt(s, Inches(0.9), Inches(4.15), Inches(11.5), Inches(0.7),
        [("AI × 不动产资产管理高级研修班  ·  2026 年 11 月首期", 16, RGBColor(0xC5, 0xD4, 0xE4), False)])
    txt(s, Inches(0.9), Inches(4.7), Inches(11.5), Inches(0.55),
        [("主办：复旦大学住房政策研究中心  ·  上海市杨浦区科技企业联合会  ·  上海市工商联房地产商会资产管理分会",
          11, MUTED, False)])

    chips = ["L1 中心名义", "品牌费书面确认", "三方分成锁定", "11 月开班"]
    x = Inches(0.9)
    for c in chips:
        rect(s, x, Inches(5.9), Inches(2.7), Inches(0.85), NAVY2, round_=True, radius=0.1)
        txt(s, x, Inches(5.9), Inches(2.7), Inches(0.85),
            [(c, 16, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += Inches(3.0)


def main():
    os.makedirs(OUT_DIR, exist_ok=True)
    cover()
    toc()
    opportunity()
    positioning()
    insight()
    audience()
    product()
    curriculum()
    faculty()
    community()
    hosts()
    fudan_brand()
    models()
    pricing()
    economics()
    ecosystem()
    timeline()
    next_steps()
    closing()
    prs.save(OUT)
    print(f"已生成 {len(prs.slides)} 页：{OUT}")


if __name__ == "__main__":
    main()
