#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成《七款中文 AI PPT Skill 横评》16:9 可编辑 PPTX。

口径：
  · 排名沿用社区锐评（彼得潘AI / 即刻），看的是「能不能做出像样的 PPT」
  · Star 数为 2026-08-23 从 GitHub 抓取的实时值
  · 仓库能力来自各项目 README / SKILL.md，社区原评单独标注

运行：python3 scripts/build_ppt_skill_overview.py
输出：output/七款中文AI-PPT-Skill横评.pptx
"""

from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ---------------------------------------------------------------- 瑞士国际主义色板
INK = RGBColor(0x11, 0x11, 0x11)
PAPER = RGBColor(0xF4, 0xF0, 0xE6)
PAPER_DEEP = RGBColor(0xE8, 0xE2, 0xD4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xE3, 0x1C, 0x24)
GRAY = RGBColor(0x6B, 0x6B, 0x6B)
GRAY_LIGHT = RGBColor(0x9A, 0x96, 0x8C)
LINE = RGBColor(0xD6, 0xD0, 0xC2)
MUTED = RGBColor(0x4A, 0x4A, 0x4A)

FONT = "微软雅黑"
PAGE_W = Inches(13.333)
PAGE_H = Inches(7.5)
FOOTER = "七款中文 AI PPT Skill 横评  ·  社区锐评版  ·  2026.08"
TOTAL = 18

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT_DIR = os.path.join(ROOT, "output")
OUT = os.path.join(OUT_DIR, "七款中文AI-PPT-Skill横评.pptx")

SKILLS = [
    {
        "rank": "01",
        "name": "ppt-master",
        "author": "hugohe",
        "stars_raw": "48,728",
        "format": "可编辑 PPTX",
        "license": "MIT",
        "lang": "Python",
        "one": "全场唯一正经做 PPT 的",
        "quote": "元素全可编辑，自带音色克隆和旁白生成，纯纯降维打击。",
        "repo": "github.com/hugohe3/ppt-master",
        "install": "npx skills add hugohe3/ppt-master",
        "use": "/ppt-master 帮我做一份产品发布 PPT，15 页，面向投资人",
        "fit": "领导汇报、客户交付、后续还要在 PowerPoint / WPS 里改。",
        "not": "只求一场炫酷网页演示、不需要办公软件兼容。",
        "points": [
            "真正的原生 PPTX：形状、图表、表格、母版、备注都可继续编辑",
            "从 PDF / Word / 网页 / Markdown 出发，先把论证理清再排版",
            "可选音色克隆与旁白，备注能转成配音甚至视频",
            "可跟随你自己的 PPT 模板；数据留在本地，不绑死某一家模型",
        ],
    },
    {
        "rank": "02",
        "name": "frontend-slides",
        "author": "张咋啦  @zarazhangrui",
        "stars_raw": "27,997",
        "format": "单文件 HTML",
        "license": "MIT",
        "lang": "JavaScript",
        "one": "主观审美最佳，完成度极高",
        "quote": "目前呈现为 HTML 格式，稍微考验一点使用者的基础。",
        "repo": "github.com/zarazhangrui/frontend-slides",
        "install": "npx skills add zarazhangrui/frontend-slides",
        "use": "/frontend-slides 帮我做一份关于 AI 行业趋势的分享",
        "fit": "技术分享、公开演讲、在线传播，追求视觉上限。",
        "not": "必须丢进 PowerPoint，给不会开浏览器的同事继续改。",
        "points": [
            "先出三张视觉预览再开工，用眼睛选风格，而不是用形容词猜",
            "12 套安全预设 + 34 套大胆模板，刻意避开「紫渐变白底」AI 味",
            "零依赖单文件 HTML；也可把现有 PPTX 抽文案后重做网页演示",
            "可部署到 Vercel 分享，也可导出 PDF 快照",
        ],
    },
    {
        "rank": "03",
        "name": "huashu-design",
        "author": "花叔  @AlchainHust",
        "stars_raw": "23,423",
        "format": "HTML + 可编辑 PPTX",
        "license": "MIT",
        "lang": "HTML",
        "one": "审美极佳，关键是能输出可编辑 PPTX",
        "quote": "不止做 PPT，是一间装进 Agent 里的设计工作室。",
        "repo": "github.com/alchaincyf/huashu-design",
        "install": "npx skills add alchaincyf/huashu-design",
        "use": "/huashu-design 帮我做一份 AI 写作工具的产品介绍 PPT",
        "fit": "品牌发布、设计要求高、还想把文件交给团队继续改。",
        "not": "只要快速出一版能念的稿，不在乎设计流程。",
        "points": [
            "任何新视觉 100% 先出三个方向初稿，选定后再批量展开",
            "内置 20 种设计哲学 + 5 维专家评审，质量有兜底",
            "高保真原型、幻灯片、动画、信息图都能做，可导出 MP4",
            "能输出可编辑 PPTX，美观和办公流转可以同时要",
        ],
    },
    {
        "rank": "04",
        "name": "guizang-ppt-skill",
        "author": "歸藏  @op7418",
        "stars_raw": "24,666",
        "format": "单文件 HTML",
        "license": "AGPL-3.0",
        "lang": "HTML",
        "one": "人上人瑞士风审美，很适合线下分享",
        "quote": "自带快捷键。电子杂志风 × 瑞士国际主义，双击浏览器就是一场体面的分享。",
        "repo": "github.com/op7418/guizang-ppt-skill",
        "install": "npx skills add op7418/guizang-ppt-skill",
        "use": "/guizang-ppt-skill 帮我做一份杂志风 PPT，主题是 AI 趋势，10 页",
        "fit": "线下分享、观点表达、个人风格演示。",
        "not": "需要团队在 PowerPoint 里来回改稿。",
        "points": [
            "Style A 电子杂志 × 电子墨水；Style B 瑞士网格 + 单一锚点色",
            "单文件 HTML，键盘 / 滚轮 / 触屏都能翻页，可全屏演示",
            "内置排练与演讲者模式，可补演讲备注和每页计划时长",
            "顺手能出配图、公众号头图、多平台封面",
        ],
    },
    {
        "rank": "05",
        "name": "html-ppt-skill",
        "author": "Lewis",
        "stars_raw": "8,023",
        "format": "单文件 HTML",
        "license": "MIT",
        "lang": "HTML",
        "one": "功能设计十分贴心",
        "quote": "自带计时器和逐字稿等实用组件，帮你把一场分享讲好。",
        "repo": "github.com/lewislulu/html-ppt-skill",
        "install": "npx skills add lewislulu/html-ppt-skill",
        "use": "/html-ppt-skill 帮我做一份产品汇报 PPT，带演讲计时器",
        "fit": "上台演讲、需要控时和讲稿、不太习惯脱稿。",
        "not": "把视觉上限当成第一目标。",
        "points": [
            "按 S 打开演讲者窗口：当前页、下一页、逐字稿、计时器四张磁吸卡片",
            "36 套主题 × 31 种版式 × 15 套整本模板 × 47 种动画",
            "纯静态 HTML/CSS/JS，零构建，中英文字体一等公民",
            "逐字稿规则很具体：提示信号而非照念，每页约 150–300 词",
        ],
    },
    {
        "rank": "06",
        "name": "baoyu-slide-deck",
        "author": "宝玉",
        "stars_raw": "25,263 *",
        "format": "图片卡片 / HTML",
        "license": "MIT",
        "lang": "TypeScript",
        "one": "NPC 风格偏可爱，以纯图片形式呈现",
        "quote": "更适合课件、培训、知识分享；整页可直接发社媒。",
        "repo": "github.com/JimLiu/baoyu-skills  ·  skills/baoyu-slide-deck",
        "install": "npx skills add JimLiu/baoyu-skills --skill baoyu-slide-deck",
        "use": "/baoyu-slide-deck tutorial.md --style corporate --slides 12",
        "fit": "课件、培训、知识产品、小红书式整页卡片。",
        "not": "正式商务汇报、需要别人改文字的场景。",
        "points": [
            "定位是「阅读和转发」而不是现场演示：自解释、可刷、可截图",
            "先出大纲和风格指令，再逐页生成位图，禁止用 SVG/HTML 冒充",
            "宝玉 Skills 生态覆盖面极广，这只是其中的 PPT 模块",
            "Star 数是整个 baoyu-skills 仓库，不是单技能",
        ],
    },
    {
        "rank": "07",
        "name": "qiaomu-anything-to-notebooklm",
        "author": "乔木",
        "stars_raw": "5,783",
        "format": "图片卡片 / 多格式",
        "license": "MIT",
        "lang": "Python",
        "one": "偏向纯图片卡片，侧重内容的初步呈现",
        "quote": "严格来说不是专业 PPT Skill，更适合做内容整理的第一步。",
        "repo": "github.com/joeseesun/qiaomu-anything-to-notebooklm",
        "install": "npx skills add joeseesun/qiaomu-anything-to-notebooklm",
        "use": "/qiaomu-anything-to-notebooklm 把这篇文章转成知识卡片",
        "fit": "把公众号、播客、PDF、网页先变成可消化的结构。",
        "not": "直接当终稿演示文件交给客户。",
        "points": [
            "多源接入：微信、网页、YouTube、PDF、播客、Office 文档",
            "经 NotebookLM 输出播客 / PPT / 思维导图 / Quiz / 闪卡",
            "适合「先看见结构」，再用上面任意一款做正式演示",
            "入门门槛低，但产物颗粒度停在内容呈现，不是设计终稿",
        ],
    },
]


# ---------------------------------------------------------------- 基础工具
def _set_run_font(run, size, color, bold=False, italic=False, name=FONT):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def add_text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.12):
    """lines: [[(text, size, color, bold), ...], ...] 或 [(text, size, color, bold)]"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass
    for m in (tf.margin_left,):
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
    first = True
    for para in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = line_spacing
        runs = para if (para and isinstance(para[0], tuple)) else [para]
        for item in runs:
            text, size, color, bold = item[0], item[1], item[2], item[3]
            italic = item[4] if len(item) > 4 else False
            r = p.add_run()
            r.text = text
            _set_run_font(r, size, color, bold, italic)
    return tb


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75,
             shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def footer(slide, idx, dark=False):
    color = RGBColor(0x88, 0x88, 0x88) if dark else GRAY_LIGHT
    add_text(slide, Inches(0.45), Inches(7.18), Inches(10.2), Inches(0.22),
             [[(FOOTER, 9, color, False)]])
    add_text(slide, Inches(11.4), Inches(7.18), Inches(1.45), Inches(0.22),
             [[(f"{idx:02d}  /  {TOTAL:02d}", 9, color, False)]],
             align=PP_ALIGN.RIGHT)


def kicker_bar(slide, kicker, title, subtitle=None):
    set_bg(slide, PAPER)
    add_rect(slide, 0, 0, PAGE_W, Inches(0.08), RED)
    add_text(slide, Inches(0.5), Inches(0.22), Inches(12.2), Inches(0.28),
             [[(kicker, 11, RED, True)]])
    add_text(slide, Inches(0.5), Inches(0.48), Inches(12.2), Inches(0.5),
             [[(title, 26, INK, True)]])
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.00), Inches(12.2), Inches(0.32),
                 [[(subtitle, 12, GRAY, False)]])


def hairline(slide, x, y, w, color=LINE):
    add_rect(slide, x, y, w, Pt(1.0), color)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---------------------------------------------------------------- 各页
def slide_cover(prs):
    s = blank(prs)
    set_bg(s, INK)
    add_rect(s, 0, 0, Inches(0.18), PAGE_H, RED)
    add_text(s, Inches(0.7), Inches(0.55), Inches(11.5), Inches(0.35),
             [[("PRESENTATION PROJECT OVERVIEW  ·  2026.08.23", 12, RED, True)]])
    add_text(s, Inches(0.7), Inches(1.55), Inches(12.0), Inches(2.2),
             [[("七款中文 AI", 42, WHITE, True)],
              [("PPT Skill 横评", 42, WHITE, True)]],
             line_spacing=1.02)
    add_rect(s, Inches(0.7), Inches(3.85), Inches(1.6), Inches(0.06), RED)
    add_text(s, Inches(0.7), Inches(4.1), Inches(11.5), Inches(1.3),
             [[("从可编辑 PPTX 到杂志风网页演示。", 18, PAPER, False)],
              [("排名沿用社区锐评，Star 数为今日 GitHub 实测。", 18, PAPER, False)]],
             line_spacing=1.2)
    add_text(s, Inches(0.7), Inches(6.55), Inches(11.5), Inches(0.4),
             [[("hugohe  ·  张咋啦  ·  花叔  ·  歸藏  ·  Lewis  ·  宝玉  ·  乔木", 13, GRAY_LIGHT, False)]])
    add_notes(s, "开场：这不是又一份「AI 做 PPT」教程，而是把社区里被反复点名的七款 Skill 摊开，讲清楚各自交什么文件、适合什么场合、什么时候不该用。")


def slide_how_to_read(prs, idx):
    s = blank(prs)
    kicker_bar(s, "01  阅读约定", "这份材料怎么用",
               "先对齐口径，再看排名。否则 Star 数字会把选型带偏。")
    cards = [
        ("排名从哪来", "社区锐评（彼得潘AI 视频 / 即刻转述），标准是「能不能做出像样的 PPT」，不是 Star 排行榜。"),
        ("Star 怎么读", "今日 GitHub 实时数。宝玉是整个 baoyu-skills 仓库；歸藏 Star 已超过花叔，但锐评名次未改。"),
        ("产物分三类", "可编辑 PPTX · 单文件 HTML 演示 · 整页图片卡片。先问「文件要交给谁」，再问「好不好看」。"),
        ("怎么选型", "没有全场第一。投资人路演、技术分享、课件卡片、内容初稿，工具应该完全不同。"),
    ]
    for i, (t, b) in enumerate(cards):
        col, row = i % 2, i // 2
        x = Inches(0.5) + col * Inches(6.35)
        y = Inches(1.55) + row * Inches(2.55)
        add_rect(s, x, y, Inches(6.1), Inches(2.35), WHITE, LINE, 1.0)
        add_rect(s, x, y, Inches(0.1), Inches(2.35), RED)
        add_text(s, x + Inches(0.35), y + Inches(0.28), Inches(5.5), Inches(0.4),
                 [[(t, 16, INK, True)]])
        add_text(s, x + Inches(0.35), y + Inches(0.85), Inches(5.5), Inches(1.25),
                 [[(b, 14, MUTED, False)]], line_spacing=1.28)
    footer(s, idx)
    add_notes(s, "提醒听众：后面每一页的「社区原评」是原文，Star 是我们今天核实过的。宝玉那一行尤其容易误读。")


def slide_toc(prs, idx):
    s = blank(prs)
    kicker_bar(s, "02  目录", "十八页，三件事")
    items = [
        ("I", "地图", "为什么是现在 · 排名一览 · 三条技术路线"),
        ("II", "七款细读", "从 ppt-master 到乔木，各一页讲清楚"),
        ("III", "怎么选", "对照矩阵 · 场景决策 · 安装速查 · 五个判断"),
    ]
    for i, (num, title, desc) in enumerate(items):
        y = Inches(1.7) + i * Inches(1.65)
        add_text(s, Inches(0.55), y, Inches(1.4), Inches(1.1),
                 [[(num, 36, RED, True)]], anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(2.2), y + Inches(0.15), Inches(10), Inches(0.45),
                 [[(title, 22, INK, True)]])
        add_text(s, Inches(2.2), y + Inches(0.65), Inches(10), Inches(0.4),
                 [[(desc, 14, GRAY, False)]])
        if i < 2:
            hairline(s, Inches(2.2), y + Inches(1.35), Inches(10.3))
    footer(s, idx)
    add_notes(s, "节奏：地图约 3 分钟，七款细读每款约 1 分钟，选型约 4 分钟。十五分钟可以讲完。")


def slide_why_now(prs, idx):
    s = blank(prs)
    kicker_bar(s, "03  背景", "2026，AI 做 PPT 终于分成了几条路",
               "以前比的是「像不像模板」；现在比的是交什么文件、能不能继续改、现场怎么讲。")
    facts = [
        ("48,728", "ppt-master 当前 Star", "独立 PPT 项目里最猛的一个，走的是 PowerPoint 原生对象模型。"),
        ("三条", "主流产物形态", "PPTX 办公流 · HTML 演示流 · 图片卡片流，已经不是同一种「PPT」。"),
        ("Skill", "跑在 Agent 里", "不是又一个在线 SaaS。装进 Claude Code / Codex / Cursor，本地出稿。"),
    ]
    for i, (n, t, d) in enumerate(facts):
        x = Inches(0.5) + i * Inches(4.2)
        add_rect(s, x, Inches(1.7), Inches(3.95), Inches(4.55), WHITE)
        add_rect(s, x, Inches(1.7), Inches(3.95), Inches(0.08), RED)
        add_text(s, x + Inches(0.3), Inches(2.1), Inches(3.35), Inches(1.1),
                 [[(n, 32, RED, True)]])
        add_text(s, x + Inches(0.3), Inches(3.3), Inches(3.35), Inches(0.55),
                 [[(t, 16, INK, True)]])
        add_text(s, x + Inches(0.3), Inches(4.0), Inches(3.35), Inches(1.8),
                 [[(d, 14, MUTED, False)]], line_spacing=1.28)
    footer(s, idx)
    add_notes(s, "点题：社区原文说「大家已经不是做同一种 PPT 了」。后面所有选型都建立在这句话上。")


def slide_ranking(prs, idx):
    s = blank(prs)
    kicker_bar(s, "04  一览", "社区锐评排名 × 今日 Star",
               "排名未按 Star 重排。歸藏当前 Star 已高于花叔；宝玉为整个技能仓库。")
    headers = ["#", "Skill", "作者", "今日 Star", "产物", "社区原评"]
    rows = [
        ["01", "ppt-master", "hugohe", "48,728", "可编辑 PPTX", "全场唯一正经做 PPT · 音色克隆"],
        ["02", "frontend-slides", "张咋啦", "27,997", "HTML", "主观审美最佳 · 完成度极高"],
        ["03", "huashu-design", "花叔", "23,423", "HTML + PPTX", "审美极佳 · 可编辑 PPTX"],
        ["04", "guizang-ppt-skill", "歸藏", "24,666", "HTML", "瑞士风 · 快捷键 · 线下分享"],
        ["05", "html-ppt-skill", "Lewis", "8,023", "HTML", "计时器 + 逐字稿"],
        ["06", "baoyu-slide-deck", "宝玉", "25,263*", "图片卡片", "NPC 可爱风 · 整页图片"],
        ["07", "qiaomu-to-notebooklm", "乔木", "5,783", "图片 / 多格式", "内容初步呈现 · 入门"],
    ]
    col_w = [Inches(0.7), Inches(2.55), Inches(1.45), Inches(1.45), Inches(2.15), Inches(4.15)]
    x0, y0, row_h = Inches(0.5), Inches(1.52), Inches(0.68)
    add_rect(s, x0, y0, Inches(12.35), Inches(0.42), INK)
    x = x0
    for i, h in enumerate(headers):
        add_text(s, x + Inches(0.08), y0, col_w[i] - Inches(0.08), Inches(0.42),
                 [[(h, 11, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[i]
    for r, row in enumerate(rows):
        y = y0 + Inches(0.42) + r * row_h
        bg = WHITE if r % 2 == 0 else PAPER_DEEP
        add_rect(s, x0, y, Inches(12.35), row_h, bg)
        x = x0
        for c, val in enumerate(row):
            color = RED if c == 0 else INK
            bold = c in (0, 1)
            add_text(s, x + Inches(0.08), y, col_w[c] - Inches(0.1), row_h,
                     [[(val, 12, color, bold)]], anchor=MSO_ANCHOR.MIDDLE)
            x += col_w[c]
    add_text(s, Inches(0.5), Inches(6.82), Inches(12.3), Inches(0.28),
             [[("* 宝玉 Star 为 JimLiu/baoyu-skills 整仓，不是 baoyu-slide-deck 单技能。", 10, GRAY, False)]])
    footer(s, idx)
    add_notes(s, "不要在这一页逐行念。指给听众：真正的分叉在「产物」这一列。")


def slide_three_routes(prs, idx):
    s = blank(prs)
    kicker_bar(s, "05  分叉", "先选路线，再选 Skill")
    routes = [
        ("A", "可编辑 PPTX", "ppt-master\nhuashu-design",
         "团队都能打开改。图表、形状、母版是真对象。动效和版式上限不如网页。",
         "领导汇报 · 客户交付 · 模板继承"),
        ("B", "单文件 HTML", "frontend-slides\nguizang-ppt-skill\nhtml-ppt-skill",
         "双击浏览器就能演示，动效和排版自由。发给不懂技术的同事有门槛。",
         "技术分享 · 线下演讲 · 在线传播"),
        ("C", "整页图片卡片", "baoyu-slide-deck\nqiaomu-anything-to-notebooklm",
         "适合刷、截、发。文字改不了，不是正式演示格式。",
         "课件 · 社媒 · 内容初稿"),
    ]
    for i, (letter, title, who, desc, scene) in enumerate(routes):
        x = Inches(0.45) + i * Inches(4.25)
        add_rect(s, x, Inches(1.55), Inches(4.05), Inches(5.15), WHITE)
        add_rect(s, x, Inches(1.55), Inches(4.05), Inches(0.9), INK)
        add_text(s, x + Inches(0.22), Inches(1.55), Inches(0.6), Inches(0.9),
                 [[(letter, 22, RED, True)]], anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.85), Inches(1.55), Inches(3.0), Inches(0.9),
                 [[(title, 16, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.25), Inches(2.6), Inches(3.55), Inches(1.15),
                 [[(who, 13, RED, True)]], line_spacing=1.25)
        add_text(s, x + Inches(0.25), Inches(3.85), Inches(3.55), Inches(1.5),
                 [[(desc, 13, MUTED, False)]], line_spacing=1.25)
        add_rect(s, x + Inches(0.25), Inches(5.5), Inches(3.55), Inches(0.9), PAPER_DEEP)
        add_text(s, x + Inches(0.4), Inches(5.5), Inches(3.25), Inches(0.9),
                 [[(scene, 12, INK, True)]], anchor=MSO_ANCHOR.MIDDLE)
    footer(s, idx)
    add_notes(s, "决策顺序必须是：文件交给谁 → 路线 A/B/C → 再在路线里挑审美和功能。")


def slide_profile(prs, idx, skill):
    s = blank(prs)
    set_bg(s, PAPER)
    add_rect(s, 0, 0, PAGE_W, Inches(0.08), RED)

    add_rect(s, 0, Inches(0.08), Inches(3.55), Inches(7.42), INK)
    add_text(s, Inches(0.28), Inches(0.28), Inches(3.05), Inches(0.7),
             [[(skill["rank"], 40, RED, True)]])
    add_text(s, Inches(0.28), Inches(1.05), Inches(3.1), Inches(1.15),
             [[(skill["name"], 16, WHITE, True)]], line_spacing=1.05)
    add_text(s, Inches(0.28), Inches(2.25), Inches(3.1), Inches(0.7),
             [[(skill["author"], 12, GRAY_LIGHT, False)]])
    stats = [
        ("今日 Star", skill["stars_raw"]),
        ("产物", skill["format"]),
        ("协议", skill["license"]),
        ("实现", skill["lang"]),
    ]
    for i, (k, v) in enumerate(stats):
        y = Inches(3.15) + i * Inches(0.72)
        add_text(s, Inches(0.28), y, Inches(3.1), Inches(0.22),
                 [[(k, 10, GRAY_LIGHT, False)]])
        add_text(s, Inches(0.28), y + Inches(0.22), Inches(3.1), Inches(0.4),
                 [[(v, 13, WHITE, True)]])

    add_text(s, Inches(3.85), Inches(0.28), Inches(8.9), Inches(0.28),
             [[("社区原评", 11, RED, True)]])
    add_text(s, Inches(3.85), Inches(0.55), Inches(8.9), Inches(0.5),
             [[(skill["one"], 20, INK, True)]])
    add_text(s, Inches(3.85), Inches(1.1), Inches(8.9), Inches(0.55),
             [[(skill["quote"], 13, MUTED, False)]], line_spacing=1.2)

    hairline(s, Inches(3.85), Inches(1.75), Inches(8.9))

    add_text(s, Inches(3.85), Inches(1.9), Inches(8.9), Inches(0.3),
             [[("它实际在做的事", 12, RED, True)]])
    y = Inches(2.25)
    for p in skill["points"]:
        add_rect(s, Inches(3.85), y + Inches(0.1), Inches(0.08), Inches(0.08), RED)
        add_text(s, Inches(4.1), y, Inches(8.6), Inches(0.48),
                 [[(p, 13, INK, False)]])
        y += Inches(0.5)

    add_rect(s, Inches(3.85), Inches(4.4), Inches(4.3), Inches(1.55), WHITE, LINE, 1.0)
    add_text(s, Inches(4.05), Inches(4.5), Inches(4.0), Inches(0.3),
             [[("适合", 12, RED, True)]])
    add_text(s, Inches(4.05), Inches(4.85), Inches(3.95), Inches(0.95),
             [[(skill["fit"], 13, INK, False)]], line_spacing=1.2)

    add_rect(s, Inches(8.35), Inches(4.4), Inches(4.4), Inches(1.55), WHITE, LINE, 1.0)
    add_text(s, Inches(8.55), Inches(4.5), Inches(4.05), Inches(0.3),
             [[("不适合", 12, RED, True)]])
    add_text(s, Inches(8.55), Inches(4.85), Inches(4.05), Inches(0.95),
             [[(skill["not"], 13, INK, False)]], line_spacing=1.2)

    add_rect(s, Inches(3.85), Inches(6.1), Inches(8.9), Inches(0.85), INK)
    add_text(s, Inches(4.05), Inches(6.15), Inches(8.5), Inches(0.28),
             [[(skill["repo"], 11, GRAY_LIGHT, False)]])
    add_text(s, Inches(4.05), Inches(6.42), Inches(8.5), Inches(0.42),
             [[(skill["install"], 12, WHITE, True)]])

    footer(s, idx)
    add_notes(s, f"{skill['name']}：{skill['one']} 使用口令：{skill['use']}")


def slide_matrix(prs, idx):
    s = blank(prs)
    kicker_bar(s, "13  对照", "七维速查：别只看审美")
    headers = ["Skill", "PPTX", "HTML", "旁白/讲稿", "审美上限", "办公流转", "现场演示"]
    data = [
        ["ppt-master", "原生可改", "—", "音色克隆", "高", "最强", "中"],
        ["frontend-slides", "可转入", "最强", "—", "最高", "弱", "强"],
        ["huashu-design", "可导出", "很强", "动画/MP4", "很高", "强", "强"],
        ["guizang-ppt-skill", "—", "很强", "备注/排练", "很高", "弱", "最强之一"],
        ["html-ppt-skill", "—", "强", "计时+逐字稿", "高", "弱", "最贴心"],
        ["baoyu-slide-deck", "可拼图", "可", "—", "可爱风", "弱", "偏阅读"],
        ["qiaomu", "初稿级", "—", "播客", "入门", "弱", "弱"],
    ]
    col_w = [Inches(2.35), Inches(1.55), Inches(1.45), Inches(1.85),
             Inches(1.45), Inches(1.55), Inches(2.15)]
    x0, y0, rh = Inches(0.5), Inches(1.5), Inches(0.68)
    add_rect(s, x0, y0, Inches(12.35), Inches(0.42), INK)
    x = x0
    for i, h in enumerate(headers):
        add_text(s, x + Inches(0.08), y0, col_w[i] - Inches(0.08), Inches(0.42),
                 [[(h, 11, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[i]
    for r, row in enumerate(data):
        y = y0 + Inches(0.42) + r * rh
        add_rect(s, x0, y, Inches(12.35), rh, WHITE if r % 2 == 0 else PAPER_DEEP)
        x = x0
        for c, val in enumerate(row):
            color = INK if c == 0 else MUTED
            add_text(s, x + Inches(0.08), y, col_w[c] - Inches(0.1), rh,
                     [[(val, 12, color, c == 0)]], anchor=MSO_ANCHOR.MIDDLE)
            x += col_w[c]
    footer(s, idx)
    add_notes(s, "「审美上限」是社区主观评价。办公流转几乎只看能不能交出别人打得开的 PPTX。")


def slide_scenes(prs, idx):
    s = blank(prs)
    kicker_bar(s, "14  选型", "不做排名，只看你在干什么")
    scenes = [
        ("投资人路演 / 领导汇报", "ppt-master", "要可继续改的真 PPTX，旁白是加分项。"),
        ("品牌发布 / 设计很挑", "huashu-design", "先选方向稿，再要 PPTX 和动画。"),
        ("技术分享 / 视觉上限", "frontend-slides", "三张预览定风格，单文件 HTML 发出去。"),
        ("线下瑞士风 / 杂志风", "guizang-ppt-skill", "快捷键翻页，双击浏览器开讲。"),
        ("上台要计时和讲稿", "html-ppt-skill", "演讲者模式 + 逐字稿，控场优先。"),
        ("课件 / 社媒整页卡片", "baoyu-slide-deck", "为阅读和转发设计，不是为投影。"),
        ("先把素材变成结构", "qiaomu-anything-to-notebooklm", "初稿和多格式转化，不是终稿。"),
    ]
    for i, (scene, skill, why) in enumerate(scenes):
        y = Inches(1.45) + i * Inches(0.72)
        add_rect(s, Inches(0.5), y, Inches(12.35), Inches(0.64), WHITE)
        add_rect(s, Inches(0.5), y, Inches(0.1), Inches(0.64), RED)
        add_text(s, Inches(0.8), y, Inches(4.4), Inches(0.64),
                 [[(scene, 14, INK, True)]], anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(5.3), y, Inches(3.3), Inches(0.64),
                 [[(skill, 13, RED, True)]], anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(8.6), y, Inches(4.05), Inches(0.64),
                 [[(why, 12, MUTED, False)]], anchor=MSO_ANCHOR.MIDDLE)
    footer(s, idx)
    add_notes(s, "Codex 团队的框架更务实：不排高低，只看场景。这一页就是那套框架的落地。")


def slide_install(prs, idx):
    s = blank(prs)
    kicker_bar(s, "15  上手", "一条命令装进 Agent",
               "前提：本机已有 Claude Code / Codex / 其他支持 Agent Skills 的工具。")
    cmds = [
        ("01  ppt-master", "npx skills add hugohe3/ppt-master"),
        ("02  frontend-slides", "npx skills add zarazhangrui/frontend-slides"),
        ("03  huashu-design", "npx skills add alchaincyf/huashu-design"),
        ("04  guizang-ppt-skill", "npx skills add op7418/guizang-ppt-skill"),
        ("05  html-ppt-skill", "npx skills add lewislulu/html-ppt-skill"),
        ("06  baoyu-slide-deck", "npx skills add JimLiu/baoyu-skills --skill baoyu-slide-deck"),
        ("07  qiaomu", "npx skills add joeseesun/qiaomu-anything-to-notebooklm"),
    ]
    for i, (label, cmd) in enumerate(cmds):
        y = Inches(1.48) + i * Inches(0.7)
        add_rect(s, Inches(0.5), y, Inches(12.35), Inches(0.62), WHITE)
        add_text(s, Inches(0.7), y, Inches(3.5), Inches(0.62),
                 [[(label, 13, RED, True)]], anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(4.3), y, Inches(8.3), Inches(0.62),
                 [[(cmd, 13, INK, False)]], anchor=MSO_ANCHOR.MIDDLE)
    footer(s, idx)
    add_notes(s, "安装命令来自 01MVP 横评与各仓库 README。具体 Agent 的技能目录可能还要手动拷 SKILL.md。")


def slide_judgements(prs, idx):
    s = blank(prs)
    kicker_bar(s, "16  判断", "五句把这盘棋下完")
    items = [
        ("01", "可编辑已经是入场券，不是奖牌。",
         "ppt-master 真正拉开差距的是原生对象模型：形状把手、数据图表、母版继承，以及还在向 PowerPoint 本身靠拢。"),
        ("02", "HTML 路线赢在现场，输在流转。",
         "frontend-slides / 歸藏 / Lewis 双击就能讲；但文件交给行政、法务、客户时，浏览器会变成门槛。"),
        ("03", "花叔走的是工作室流程，不是模板填充。",
         "三方向初稿是硬门。五页以上先定视觉语言，再批量展开——慢在前面，稳在后面。"),
        ("04", "宝玉和乔木不该被放进「正经 PPT」赛道硬比。",
         "一个为阅读转发做整页图，一个把多源内容先变成结构。第六、第七名不是差，是需求不同。"),
        ("05", "选工具之前先写一句：这份东西第二天要交给谁。",
         "交给领导改 → PPTX。自己上台 → HTML。发群发社媒 → 图片。这句话比任何排名都管用。"),
    ]
    for i, (num, title, body) in enumerate(items):
        y = Inches(1.42) + i * Inches(1.05)
        add_text(s, Inches(0.5), y, Inches(0.7), Inches(0.9),
                 [[(num, 16, RED, True)]])
        add_text(s, Inches(1.25), y, Inches(11.5), Inches(0.35),
                 [[(title, 15, INK, True)]])
        add_text(s, Inches(1.25), y + Inches(0.38), Inches(11.5), Inches(0.55),
                 [[(body, 13, MUTED, False)]])
    footer(s, idx)
    add_notes(s, "这五句是整份材料的结论页。如果只能留一页，留这一页。")


def slide_close(prs, idx):
    s = blank(prs)
    set_bg(s, INK)
    add_rect(s, 0, 0, Inches(0.18), PAGE_H, RED)
    add_text(s, Inches(0.7), Inches(1.7), Inches(12), Inches(0.4),
             [[("没有最好的 PPT Skill", 14, RED, True)]])
    add_text(s, Inches(0.7), Inches(2.2), Inches(12), Inches(1.8),
             [[("只有最匹配", 36, WHITE, True)],
              [("你当前场景的那一款", 36, WHITE, True)]],
             line_spacing=1.05)
    add_rect(s, Inches(0.7), Inches(4.3), Inches(1.6), Inches(0.06), RED)
    add_text(s, Inches(0.7), Inches(4.6), Inches(11.5), Inches(1.2),
             [[("路演交 ppt-master。上台交 frontend-slides 或歸藏。", 16, PAPER, False)],
              [("要设计流程交花叔。要讲稿交 Lewis。要卡片交宝玉。要结构交乔木。", 16, PAPER, False)]],
             line_spacing=1.25)
    add_text(s, Inches(0.7), Inches(6.55), Inches(11.5), Inches(0.35),
             [[("资料截止 2026-08-23  ·  仓库链接见各细读页  ·  可在 PowerPoint / WPS 中继续编辑", 12, GRAY_LIGHT, False)]])
    add_notes(s, "收束：把七个名字还回给场景。问听众自己下一次要交的文件是什么。")


def build():
    os.makedirs(OUT_DIR, exist_ok=True)
    prs = Presentation()
    prs.slide_width = PAGE_W
    prs.slide_height = PAGE_H

    slide_cover(prs)
    slide_how_to_read(prs, 2)
    slide_toc(prs, 3)
    slide_why_now(prs, 4)
    slide_ranking(prs, 5)
    slide_three_routes(prs, 6)
    for i, skill in enumerate(SKILLS):
        slide_profile(prs, 7 + i, skill)
    slide_matrix(prs, 14)
    slide_scenes(prs, 15)
    slide_install(prs, 16)
    slide_judgements(prs, 17)
    slide_close(prs, 18)

    actual = len(prs.slides)
    if actual != TOTAL:
        print(f"警告：预期 {TOTAL} 页，实际 {actual} 页")

    prs.save(OUT)
    size_kb = os.path.getsize(OUT) / 1024
    print(f"已生成：{OUT}")
    print(f"页数：{actual}    大小：{size_kb:.1f} KB")
    return OUT, actual


if __name__ == "__main__":
    build()
