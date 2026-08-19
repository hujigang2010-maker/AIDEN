#!/usr/bin/env python3
"""生成《WorkBuddy 银行引荐判断》决策简报 PPT（9页）。"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from office_compat import apply_pptx_compat


NAVY = RGBColor(0x0B, 0x2F, 0x5B)
NAVY_DARK = RGBColor(0x06, 0x1E, 0x3C)
GOLD = RGBColor(0xC4, 0xA3, 0x5A)
LIGHT = RGBColor(0xF3, 0xF6, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1F, 0x2A, 0x37)
GRAY = RGBColor(0x5C, 0x6B, 0x7A)
SOFT = RGBColor(0xE8, 0xEE, 0xF5)
GREEN = RGBColor(0x2F, 0x6B, 0x4F)
RED = RGBColor(0xA6, 0x3D, 0x2F)
AMBER = RGBColor(0xB5, 0x7A, 0x2A)

FONT = "Microsoft YaHei"
FOOTER_TEXT = "WorkBuddy 银行引荐判断 · 内部材料 · 2026-08-16"
TOTAL = 9


def add_rect(slide, x, y, w, h, fill, *, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_round(slide, x, y, w, h, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    shp.adjustments[0] = 0.08
    return shp


def set_run(run, text, *, size=16, bold=False, color=DARK, font=FONT):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    size=16,
    bold=False,
    color=DARK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font=FONT,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        set_run(p.add_run(), line, size=size, bold=bold, color=color, font=font)
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=14, color=DARK, bullet_color=GOLD, spacing=1.2):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = spacing
        set_run(p.add_run(), "●  ", size=size, color=bullet_color)
        set_run(p.add_run(), item, size=size, color=color)
    return tb


def header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SW, Inches(0.95), NAVY)
    add_rect(slide, 0, Inches(0.95), SW, Inches(0.05), GOLD)
    add_text(slide, Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.42), title, size=22, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.58), Inches(12.3), Inches(0.3), subtitle, size=12, color=LIGHT)


def footer(slide, page):
    add_text(slide, Inches(0.5), Inches(7.12), Inches(10.5), Inches(0.28), FOOTER_TEXT, size=10, color=GRAY)
    add_text(
        slide,
        Inches(11.5),
        Inches(7.12),
        Inches(1.3),
        Inches(0.28),
        f"{page} / {TOTAL}",
        size=10,
        color=GRAY,
        align=PP_ALIGN.RIGHT,
    )


def build_ppt(output_path: Path) -> None:
    global SW, SH
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    # 1 封面
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, NAVY_DARK)
    add_rect(s, 0, 0, Inches(0.18), SH, GOLD)
    add_text(s, Inches(0.8), Inches(1.15), Inches(11.5), Inches(0.35), "引荐人内部判断 · 不对外", size=14, color=GOLD)
    add_text(
        s,
        Inches(0.8),
        Inches(1.65),
        Inches(11.7),
        Inches(1.35),
        "WorkBuddy 银行引荐判断",
        size=36,
        bold=True,
        color=WHITE,
    )
    add_text(
        s,
        Inches(0.8),
        Inches(3.05),
        Inches(11.7),
        Inches(0.55),
        "要不要现在引荐中国银行、上海银行，以及下周五怎么见杨行长",
        size=16,
        color=LIGHT,
    )
    add_rect(s, Inches(0.8), Inches(3.75), Inches(2.2), Inches(0.05), GOLD)
    add_round(s, Inches(0.8), Inches(4.05), Inches(11.5), Inches(1.55), NAVY)
    add_text(
        s,
        Inches(1.05),
        Inches(4.2),
        Inches(11.05),
        Inches(1.25),
        "结论：可以引荐，但不要本周并行约见。\n"
        "先把 8 月 21 日杨行长这场走实；上海银行做第二主线预热，中国银行只做探索。\n"
        "不要用新银行催腾讯——他们若是方案没齐，多见面会更亏。",
        size=15,
        color=WHITE,
    )
    add_text(
        s,
        Inches(0.8),
        Inches(5.85),
        Inches(11.5),
        Inches(0.35),
        "已知进度：泰隆已对接　东口支行行长已见　下周五见上海分行杨行长　腾讯思考偏长",
        size=13,
        color=GOLD,
    )
    add_text(s, Inches(0.8), Inches(6.35), Inches(11.5), Inches(0.3), "2026年8月16日", size=12, color=GRAY)

    # 2 当前态势
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "一、当前态势", "泰隆已经走到分行行长，腾讯还在想，辅线不要抢戏")
    steps = [
        ("1", "已对接泰隆", "WorkBuddy 进银行渠道\n积分 / 开户赠礼 / 员工提效"),
        ("2", "东口支行已见", "一线行长验证过话题\n不是空对空介绍"),
        ("3", "周五见杨行长", "8月21日上海分行\n从网点验证跳到立项"),
        ("?", "腾讯想得久", "方案、合规或授权\n可能还没锁"),
    ]
    for i, (num, title, body) in enumerate(steps):
        x = Inches(0.4 + i * 3.2)
        bg = NAVY if i < 3 else SOFT
        tc = GOLD if i < 3 else RED
        bc = LIGHT if i < 3 else DARK
        add_round(s, x, Inches(1.25), Inches(3.0), Inches(2.55), bg)
        add_text(s, x + Inches(0.2), Inches(1.4), Inches(2.6), Inches(0.4), num, size=20, bold=True, color=tc)
        add_text(s, x + Inches(0.2), Inches(1.85), Inches(2.6), Inches(0.4), title, size=16, bold=True, color=WHITE if i < 3 else DARK)
        add_text(s, x + Inches(0.2), Inches(2.35), Inches(2.6), Inches(1.2), body, size=13, color=bc)
    add_round(s, Inches(0.4), Inches(4.05), Inches(12.5), Inches(2.7), SOFT)
    add_text(s, Inches(0.7), Inches(4.2), Inches(12.0), Inches(0.4), "为什么杨行长会比再引荐两家更值钱", size=16, bold=True, color=NAVY)
    add_bullets(
        s,
        Inches(0.7),
        Inches(4.7),
        Inches(12.0),
        Inches(1.85),
        [
            "支行行长确认“一线愿不愿意推”；分行行长才能确认“上海要不要立项”。",
            "泰隆已付出洽谈、见面、材料的沉没成本，周五是兑现窗口。中行和上海银行还在“有没有人愿意见”。",
            "引荐人的价值不在同一周堆局，而在让腾讯把能说的说清、把要的决策要回。",
        ],
        size=14,
    )
    footer(s, 2)

    # 3 两种病因
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "二、先问腾讯为什么慢", "两种病因对策相反，问清之前默认保护主线")
    add_round(s, Inches(0.4), Inches(1.25), Inches(6.15), Inches(5.5), SOFT)
    add_rect(s, Inches(0.4), Inches(1.25), Inches(0.12), Inches(5.5), GREEN)
    add_text(s, Inches(0.75), Inches(1.45), Inches(5.6), Inches(0.4), "病因 A　缺银行侧真实需求", size=18, bold=True, color=GREEN)
    add_text(s, Inches(0.75), Inches(1.95), Inches(5.6), Inches(0.7), "内部想推，但还在问：银行到底买不买账？除了泰隆还有没有人愿意谈？", size=14, color=DARK)
    add_text(s, Inches(0.75), Inches(2.7), Inches(5.6), Inches(0.35), "这时加线有用", size=13, bold=True, color=GREEN)
    add_bullets(
        s,
        Inches(0.75),
        Inches(3.15),
        Inches(5.55),
        Inches(3.2),
        [
            "证明需求不只泰隆一家",
            "降低“只押一单”的心理负担",
            "会后正式引荐：先上海银行，再中行",
            "本周仍只预热，不占会前精力",
        ],
        size=14,
        bullet_color=GREEN,
    )

    add_round(s, Inches(6.75), Inches(1.25), Inches(6.15), Inches(5.5), SOFT)
    add_rect(s, Inches(6.75), Inches(1.25), Inches(0.12), Inches(5.5), RED)
    add_text(s, Inches(7.1), Inches(1.45), Inches(5.6), Inches(0.4), "病因 B　内部包装没好", size=18, bold=True, color=RED)
    add_text(s, Inches(7.1), Inches(1.95), Inches(5.6), Inches(0.7), "价格、发票、合规、品牌授权、谁来拍板还没对齐，对银行只能“再沟通一下”。", size=14, color=DARK)
    add_text(s, Inches(7.1), Inches(2.7), Inches(5.6), Inches(0.35), "这时加线有害", size=13, bold=True, color=RED)
    add_bullets(
        s,
        Inches(7.1),
        Inches(3.15),
        Inches(5.55),
        Inches(3.2),
        [
            "大行问三句就会露出缺口",
            "第一印象变成“没准备好”",
            "本周只收紧杨行长会的可说边界",
            "问清之前，默认按病因 B 处理",
        ],
        size=14,
        bullet_color=RED,
    )
    footer(s, 3)

    # 4 三家银行对照
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "三、三家银行怎么排", "能级不同，不能同一周拉上台")
    cards = [
        (NAVY, GOLD, WHITE, LIGHT, "泰隆银行", "本周唯一主线", "城商行 · 中小微最贴\n东口支行已验证\n周五见杨行长要方向\n要牵头人、试点、下次谁来\n不要把 1500 万当承诺"),
        (SOFT, NAVY, DARK, DARK, "上海银行", "最合适的第二主线", "总部在上海，决策半径短\n叙事可复用泰隆那套\n本周只探询人选和条线\n会后视信号再正式引荐\n先分行部门或支行一把手"),
        (SOFT, AMBER, DARK, DARK, "中国银行", "只做探索，不当对手盘", "品牌重、流程长、合规严\n方案不标准就容易无下文\n必须落到具体支行/条线\n本周最多问愿不愿先聊\n不当催单工具"),
    ]
    for i, (bg, title_c, body_c, sub_c, title, tag, body) in enumerate(cards):
        x = Inches(0.4 + i * 4.25)
        add_round(s, x, Inches(1.25), Inches(4.05), Inches(5.5), bg)
        add_text(s, x + Inches(0.25), Inches(1.45), Inches(3.55), Inches(0.4), title, size=20, bold=True, color=title_c if bg != NAVY else GOLD)
        add_text(s, x + Inches(0.25), Inches(1.9), Inches(3.55), Inches(0.4), tag, size=13, bold=True, color=sub_c)
        add_text(s, x + Inches(0.25), Inches(2.5), Inches(3.55), Inches(3.9), body, size=14, color=body_c)
    footer(s, 4)

    # 5 收益与风险
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "四、现在引荐：收益对风险", "方向对，时间不对")
    add_round(s, Inches(0.4), Inches(1.25), Inches(6.15), Inches(5.5), SOFT)
    add_text(s, Inches(0.7), Inches(1.45), Inches(5.6), Inches(0.4), "为什么会想引荐", size=18, bold=True, color=NAVY)
    add_bullets(
        s,
        Inches(0.7),
        Inches(2.05),
        Inches(5.6),
        Inches(4.4),
        [
            "降低只押泰隆一家的风险",
            "给腾讯外部节奏，从“做不做”变成“先做谁”",
            "验证模式能不能复制到第二家银行",
            "引荐人持续输送资源，不绑死在一单",
        ],
        size=15,
        bullet_color=GREEN,
        spacing=1.45,
    )
    add_round(s, Inches(6.75), Inches(1.25), Inches(6.15), Inches(5.5), SOFT)
    add_text(s, Inches(7.05), Inches(1.45), Inches(5.6), Inches(0.4), "为什么本周不宜并行约见", size=18, bold=True, color=RED)
    add_bullets(
        s,
        Inches(7.05),
        Inches(2.05),
        Inches(5.6),
        Inches(4.4),
        [
            "冲主线：周五会仍可能“再想想”",
            "大行问三句，腾讯若没准备更亏",
            "泰隆若听说还在见别人，不愿拍板",
            "把引荐用成催单，关系会变味",
        ],
        size=15,
        bullet_color=RED,
        spacing=1.45,
    )
    footer(s, 5)

    # 6 打法
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "五、建议打法", "引荐 ≠ 本周约见。三档动作分开")
    rows = [
        ("本周", "8/16—21", NAVY, GOLD, "只做杨行长这一场正式会。中行、上海银行进预热名单，最多一条探询。问清腾讯卡点。明确告诉腾讯：本周不安排第二场。"),
        ("会后 3 天", "8/24—25", GREEN, WHITE, "按信号升级或按住。杨行长给方向：辅线继续预热。礼貌但空：上海银行升为正式引荐。腾讯答不利落：停新引荐，先补说法。"),
        ("不要做", "红线", RED, WHITE, "会前拉腾讯见中行或上海银行；对两家新银行报 1500 万；对泰隆说“同时还在谈别人”；把引荐当催单。"),
    ]
    for i, (tag, when, bar, tag_c, body) in enumerate(rows):
        y = Inches(1.25 + i * 1.8)
        add_round(s, Inches(0.4), y, Inches(12.5), Inches(1.65), SOFT)
        add_rect(s, Inches(0.4), y, Inches(0.14), Inches(1.65), bar)
        add_text(s, Inches(0.8), y + Inches(0.18), Inches(2.2), Inches(0.4), tag, size=18, bold=True, color=bar)
        add_text(s, Inches(0.8), y + Inches(0.62), Inches(2.2), Inches(0.35), when, size=13, color=GRAY)
        add_text(s, Inches(3.2), y + Inches(0.25), Inches(9.4), Inches(1.2), body, size=15, color=DARK)
    footer(s, 6)

    # 7 杨行长约见
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "六、下周五杨行长约见", "目标不是签单，是把支行验证变成分行方向")
    goals = [
        ("1", "原则同意", "作为中小微增值服务\n继续论证"),
        ("2", "指定牵头", "部门或人落地\n避免会后再研究"),
        ("3", "划定试点", "部分网点/客群\n先小后大"),
        ("4", "约下次", "商务/科技/合规\n谁来、何时"),
    ]
    for i, (num, title, body) in enumerate(goals):
        x = Inches(0.4 + i * 3.2)
        add_round(s, x, Inches(1.2), Inches(3.05), Inches(1.7), NAVY if i == 0 else SOFT)
        c1 = GOLD if i == 0 else NAVY
        c2 = WHITE if i == 0 else DARK
        add_text(s, x + Inches(0.18), Inches(1.3), Inches(2.7), Inches(0.35), f"{num}  {title}", size=15, bold=True, color=c1)
        add_text(s, x + Inches(0.18), Inches(1.7), Inches(2.7), Inches(1.05), body, size=13, color=c2)

    add_round(s, Inches(0.4), Inches(3.1), Inches(6.15), Inches(3.6), SOFT)
    add_text(s, Inches(0.65), Inches(3.25), Inches(5.7), Inches(0.35), "可说", size=16, bold=True, color=GREEN)
    add_bullets(
        s,
        Inches(0.65),
        Inches(3.7),
        Inches(5.7),
        Inches(2.8),
        [
            "东口支行已交流，一线认为有差异化权益",
            "主方向：积分兑换、开户/拜访赠 AI",
            "员工应用、联名卡可研究，不当场拍",
            "先试点后扩量；不接核心系统",
        ],
        size=13,
        bullet_color=GREEN,
    )
    add_round(s, Inches(6.75), Inches(3.1), Inches(6.15), Inches(3.6), SOFT)
    add_text(s, Inches(7.0), Inches(3.25), Inches(5.7), Inches(0.35), "不可说", size=16, bold=True, color=RED)
    add_bullets(
        s,
        Inches(7.0),
        Inches(3.7),
        Inches(5.7),
        Inches(2.8),
        [
            "1500 万不是已定预算或报价承诺",
            "不说行业首创、腾讯联合产品",
            "不提中行、上海银行，不用别家催",
            "不请行长当场定采购金额",
        ],
        size=13,
        bullet_color=RED,
    )
    footer(s, 7)

    # 8 口径
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "七、给三方的口径", "同一周，四套话不能串")
    quotes = [
        ("对腾讯", "泰隆周五见杨行长是本周唯一正式场。中行和上海银行我可以预热人选，不占你们会前时间。你们若卡在方案和合规，我先不把大行带上场。"),
        ("对上海银行", "腾讯 WorkBuddy 在看银行渠道，上海已有一家城商行在交流。想先了解你们公司金融或零售权益条线是否方便沟通。不是采购谈判。"),
        ("对中国银行", "想介绍一个企业 AI 工具，看是否适合做客户增值或网点补充。先找上海具体支行或普惠/公司金融条线轻量交流，不谈总行级合作。"),
        ("对泰隆", "东口支行已经沟通过。今天请杨行长看两件事：值不值得在上海试点；如果值得，谁来牵头。金额和系统对接都可以后置。"),
    ]
    for i, (who, text) in enumerate(quotes):
        y = Inches(1.2 + i * 1.4)
        add_round(s, Inches(0.4), y, Inches(12.5), Inches(1.28), SOFT)
        add_rect(s, Inches(0.4), y, Inches(2.15), Inches(1.28), NAVY)
        add_text(s, Inches(0.5), y, Inches(1.95), Inches(1.28), who, size=14, bold=True, color=GOLD, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        add_text(s, Inches(2.75), y + Inches(0.18), Inches(9.9), Inches(0.95), text, size=13, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 8)

    # 9 决策树 + 行动
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "八、会后决策树与十天动作", "用杨行长会的结果，而不是用焦虑，决定要不要引荐")
    tree = [
        (GREEN, "原则同意 + 有牵头人", "泰隆继续深挖\n上海银行预热即可\n中行暂缓"),
        (AMBER, "正向但要内部再议", "给泰隆 1—2 周\n同时约上海银行\n做探索性沟通"),
        (RED, "礼貌、无牵头、无下次", "正式引荐上海银行\n中行只做一人探询"),
        (NAVY, "腾讯会上答不利落", "停止一切新引荐\n先补一页标准说法"),
    ]
    for i, (color, title, body) in enumerate(tree):
        x = Inches(0.4 + i * 3.2)
        add_round(s, x, Inches(1.2), Inches(3.05), Inches(2.35), SOFT)
        add_rect(s, x, Inches(1.2), Inches(3.05), Inches(0.1), color)
        add_text(s, x + Inches(0.15), Inches(1.4), Inches(2.75), Inches(0.7), title, size=13, bold=True, color=color)
        add_text(s, x + Inches(0.15), Inches(2.15), Inches(2.75), Inches(1.2), body, size=13, color=DARK)

    add_round(s, Inches(0.4), Inches(3.75), Inches(12.5), Inches(2.95), NAVY)
    add_text(s, Inches(0.7), Inches(3.95), Inches(12.0), Inches(0.35), "口头一句", size=14, bold=True, color=GOLD)
    add_text(
        s,
        Inches(0.7),
        Inches(4.4),
        Inches(12.0),
        Inches(1.95),
        "本周只打泰隆杨行长；中行和上海银行先预热、会后再引荐。\n"
        "两家里先上海银行，中国银行不当催单工具。\n"
        "腾讯想得久，先问清是缺样本还是缺方案，再决定加不加线。",
        size=16,
        color=WHITE,
    )
    footer(s, 9)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    apply_pptx_compat(prs)
    prs.save(str(output_path))
    print(f"已生成：{output_path}")


if __name__ == "__main__":
    out = (
        Path(__file__).resolve().parents[1]
        / "deliverables"
        / "WorkBuddy银行引荐判断_杨行长约见简报_20260816.pptx"
    )
    build_ppt(out)
