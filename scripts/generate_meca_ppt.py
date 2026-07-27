# -*- coding: utf-8 -*-
"""生成《中东华人共创共享联盟 · 战略与落地执行》PPT。

基于战略落地执行方案 V2.0，重点融合：
1. 主办方情况
2. 未来会员定位与价格
3. 会员生态定位
4. 社群搭建（微信 / Facebook / 海外社群）
5. 大学与政府背书辅助落地建议
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ----------------------------------------------------------------------------
# 主题：中东深蓝 + 沙金（避免紫白渐变 / 奶油赤陶默认风）
# ----------------------------------------------------------------------------
INK = RGBColor(0x0B, 0x1F, 0x2A)       # 主背景
INK_2 = RGBColor(0x12, 0x2F, 0x3D)     # 次级背景
TEAL = RGBColor(0x1A, 0x4A, 0x5C)      # 卡片底
TEAL_LT = RGBColor(0x2E, 0x6B, 0x7A)   # 浅青绿
SAND = RGBColor(0xC9, 0xA2, 0x6B)      # 沙金强调
SAND_LT = RGBColor(0xE8, 0xD0, 0xA0)   # 浅沙金
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF = RGBColor(0xE8, 0xEE, 0xF1)
MUTED = RGBColor(0x9A, 0xB4, 0xBE)
ACCENT = RGBColor(0xD4, 0x7A, 0x4A)    # 暖铜点缀
OK = RGBColor(0x5A, 0xA8, 0x7A)
WARN = RGBColor(0xC9, 0x8A, 0x3C)

FONT = "微软雅黑"
SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]
PAGE = 0


def _font(run, size, color, bold=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = FONT
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT)


def slide(bg=INK):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid()
    r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def rect(s, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE, line=None):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line:
        sp.line.color.rgb = line
        sp.line.width = Pt(1.25)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def rrect(s, x, y, w, h, color, line=None, adj=0.08):
    sp = rect(s, x, y, w, h, color, MSO_SHAPE.ROUNDED_RECTANGLE, line)
    try:
        sp.adjustments[0] = adj
    except Exception:
        pass
    return sp


def txt(s, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, size, color, bold[, space_after])"""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    for i, spec in enumerate(lines):
        text, size, color, bold = spec[0], spec[1], spec[2], spec[3]
        space_after = spec[4] if len(spec) > 4 else 4
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = text
        _font(run, size, color, bold)
    return tb


def header(s, title, subtitle=None, tag=None):
    rect(s, Inches(0), Inches(0), Inches(0.16), SH, SAND)
    txt(s, Inches(0.5), Inches(0.28), Inches(10.8), Inches(0.55),
        [(title, 26, WHITE, True)])
    if subtitle:
        txt(s, Inches(0.52), Inches(0.82), Inches(10.8), Inches(0.35),
            [(subtitle, 13, SAND_LT, False)])
    if tag:
        chip = rrect(s, Inches(11.15), Inches(0.32), Inches(1.85), Inches(0.42), TEAL, SAND)
        tf = chip.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = tag
        _font(r, 11, SAND_LT, True)
    rect(s, Inches(0.5), Inches(1.25), Inches(12.3), Pt(1.5), SAND)


def footer(s, n):
    txt(s, Inches(0.5), Inches(7.08), Inches(10), Inches(0.3),
        [("中东华人共创共享联盟 MECA · 战略与落地执行 V2.0", 9, MUTED, False)])
    txt(s, Inches(11.8), Inches(7.08), Inches(1.2), Inches(0.3),
        [(str(n), 9, MUTED, False)], align=PP_ALIGN.RIGHT)


def new(title=None, subtitle=None, tag=None):
    global PAGE
    PAGE += 1
    s = slide()
    if title:
        header(s, title, subtitle, tag)
    footer(s, PAGE)
    return s


def card(s, x, y, w, h, title, bullets, accent=SAND, body_size=12):
    rrect(s, x, y, w, h, TEAL)
    rect(s, x, y, Inches(0.1), h, accent)
    txt(s, x + Inches(0.22), y + Inches(0.14), w - Inches(0.4), Inches(0.4),
        [(title, 14, SAND_LT, True)])
    lines = [("• " + b, body_size, OFF, False, 3) for b in bullets]
    txt(s, x + Inches(0.22), y + Inches(0.58), w - Inches(0.4), h - Inches(0.7), lines)


def draw_table(s, x0, y0, col_widths_in, row_h_in, headers, rows,
               header_bg=SAND, header_fg=INK, zebra=(TEAL, INK_2),
               first_col_color=SAND_LT, body_size=12, header_size=12,
               gap_in=0.06):
    """按英寸绘制表格，避免 float 与 Inches 混算导致负宽度叠字。"""
    widths = [Inches(w) for w in col_widths_in]
    xs = [x0]
    for w in widths[:-1]:
        xs.append(xs[-1] + w)
    gap = Inches(gap_in)
    row_h = Inches(row_h_in)
    # 表头
    for j, h in enumerate(headers):
        rrect(s, xs[j], y0, widths[j] - gap, row_h, header_bg)
        txt(s, xs[j] + Inches(0.08), y0, widths[j] - gap - Inches(0.12), row_h,
            [(h, header_size, header_fg, True)],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 数据行
    for i, row in enumerate(rows):
        y = y0 + row_h + Inches(0.06) + i * (row_h + Inches(0.05))
        bg = zebra[0] if i % 2 == 0 else zebra[1]
        for j, val in enumerate(row):
            rrect(s, xs[j], y, widths[j] - gap, row_h, bg)
            color = first_col_color if j == 0 else WHITE
            txt(s, xs[j] + Inches(0.1), y, widths[j] - gap - Inches(0.16), row_h,
                [(str(val), body_size, color, j == 0)],
                align=PP_ALIGN.CENTER if j else PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.MIDDLE)


# ============================================================================
# 1 封面
# ============================================================================
s = new()
rect(s, Inches(8.2), 0, Inches(5.15), SH, INK_2)
rect(s, 0, Inches(0), SW, Pt(5), SAND)
rect(s, Inches(0.7), Inches(1.35), Inches(0.16), Inches(1.35), SAND)
txt(s, Inches(1.05), Inches(1.25), Inches(7), Inches(0.45),
    [("MECA · 战略落地执行方案 V2.0", 14, SAND_LT, False)])
txt(s, Inches(1.0), Inches(1.85), Inches(7.2), Inches(2.2),
    [("中东华人共创共享联盟", 40, WHITE, True, 10),
     ("1000 位创始会员计划", 26, SAND, True, 8),
     ("暨迪拜华人生活生态平台", 22, OFF, False)])
txt(s, Inches(1.05), Inches(4.55), Inches(7), Inches(1.0),
    [("以房产为入口 · 以会员为纽带 · 以生态为长期资产", 14, MUTED, False, 6),
     ("主办方 · 会员定价 · 生态定位 · 海外社群 · 背书落地", 13, MUTED, False)])
for i, t in enumerate(["可执行", "可考核", "可复制"]):
    x = Inches(1.05 + i * 1.9)
    chip = rrect(s, x, Inches(5.85), Inches(1.7), Inches(0.48), TEAL, SAND)
    tf = chip.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = t
    _font(r, 13, SAND_LT, True)
txt(s, Inches(8.55), Inches(2.2), Inches(4.3), Inches(3.5),
    [("核心交付", 14, SAND, True, 10),
     ("① 主办方与治理架构", 14, WHITE, False, 8),
     ("② 会员定位与价格体系", 14, WHITE, False, 8),
     ("③ 会员生态与合作商网络", 14, WHITE, False, 8),
     ("④ 微信 / Facebook 等社群搭建", 14, WHITE, False, 8),
     ("⑤ 大学与政府背书路径建议", 14, WHITE, False, 8),
     ("⑥ 90 天启动与财务 KPI", 14, WHITE, False, 8)])
PAGE = 1
footer(s, 1)

# ============================================================================
# 2 目录
# ============================================================================
s = new("目录", "从愿景到施工图：按模块推进", "目录")
items = [
    ("01", "问题与定位", "痛点、边界、一句话定位"),
    ("02", "主办方与背书", "谁来办、要不要政府/大学"),
    ("03", "会员产品与价格", "分层定位、年费、生命周期"),
    ("04", "会员生态定位", "七大权益 + 六大中心"),
    ("05", "社群搭建体系", "微信 / FB / 海外矩阵"),
    ("06", "商业与落地", "收入结构、财务、90 天、KPI"),
]
for i, (no, title, desc) in enumerate(items):
    col, row = i % 3, i // 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.7 + row * 2.35)
    rrect(s, x, y, Inches(3.95), Inches(2.05), TEAL)
    txt(s, x + Inches(0.3), y + Inches(0.35), Inches(3.3), Inches(0.5),
        [(no, 22, SAND, True)])
    txt(s, x + Inches(0.3), y + Inches(0.95), Inches(3.3), Inches(0.8),
        [(title, 18, WHITE, True, 4), (desc, 12, MUTED, False)])

# ============================================================================
# 3 痛点与边界
# ============================================================================
s = new("前言：六大现实痛点", "先定义问题，再谈平台价值", "问题")
pains = [
    ("可信资源缺失", "无统一、长期稳定入口", "创始会员 + 实名认证"),
    ("信息高度分散", "散落数十个微信群", "数字化平台 + 企业黄页"),
    ("熟人介绍低效", "规模小、不可复制", "撮合机制 + 需求匹配"),
    ("会员体系脆弱", "社群生命周期短", "权益 + 履约 + 续费闭环"),
    ("生活保障缺失", "住/医/教/法无标准", "七大权益 + SLA"),
    ("生态连接断裂", "人房资源资本割裂", "六大中心 + 数据沉淀"),
]
for i, (a, b, c) in enumerate(pains):
    col, row = i % 3, i // 3
    x = Inches(0.45 + col * 4.25)
    y = Inches(1.55 + row * 2.5)
    rrect(s, x, y, Inches(4.05), Inches(2.25), TEAL)
    txt(s, x + Inches(0.25), y + Inches(0.25), Inches(3.5), Inches(1.7),
        [(f"0{i+1}  {a}", 15, SAND_LT, True, 8),
         (b, 12, OFF, False, 6),
         (f"解法：{c}", 12, SAND, False)])

# ============================================================================
# 4 定位
# ============================================================================
s = new("项目定位", "中东华人全生命周期生态平台", "定位")
rrect(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(1.35), TEAL)
txt(s, Inches(0.8), Inches(1.75), Inches(11.7), Inches(1.0),
    [("一句话定位", 12, SAND, True, 4),
     ("以房产为入口、以创始会员身份为纽带、以标准化生活与商务服务为交付、"
      "以数字平台沉淀数据与信任的中东华人全生命周期生态平台。", 15, WHITE, False)])
missions = [
    ("愿景", "成为华人在中东长期发展的基础设施：资源 / 生活 / 商业 / 价值共享"),
    ("使命", "有房 · 有圈子 · 有资源 · 有事业 · 有未来"),
    ("价值观", "共创 · 共建 · 共享 · 共赢 · 长期主义 · 诚信透明 · 利他合作"),
]
for i, (t, d) in enumerate(missions):
    y = Inches(3.15 + i * 1.15)
    rrect(s, Inches(0.5), y, Inches(12.3), Inches(1.0), INK_2, TEAL_LT)
    txt(s, Inches(0.8), y + Inches(0.2), Inches(2.0), Inches(0.6),
        [(t, 16, SAND, True)], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(3.0), y + Inches(0.2), Inches(9.4), Inches(0.6),
        [(d, 14, OFF, False)], anchor=MSO_ANCHOR.MIDDLE)

# ============================================================================
# 5 主办方情况（重点）
# ============================================================================
s = new("主办方情况：谁来办、怎么办", "实体先行 · 权责清晰 · 跨境双主体协同", "主办方")
card(s, Inches(0.45), Inches(1.5), Inches(4.05), Inches(5.2),
     "建议主办架构",
     ["迪拜侧运营主体：DED / 自贸区公司（主经营）",
      "房产交付：RERA 持牌经纪实体（自有或独家合作）",
      "中国侧协作主体：上海注册公司或合作机构（获客/内容/政企对接）",
      "品牌主体：MECA 中英文商标双注册",
      "治理层：理事会 9–12 人（发起方+理事单位+创始会员代表）",
      "执行层：CEO/秘书长 + 五部门（房产/会员/生态/市场/职能）"],
     ACCENT)
card(s, Inches(4.7), Inches(1.5), Inches(4.05), Inches(5.2),
     "主办方必须先做的事",
     ["完成法律意见书与牌照路径（P0）",
      "明确“联盟”名称使用与活动许可",
      "厘清：谁收会员费、谁签合作商、谁承担 SLA",
      "房产优惠/返佣全部写入书面合同",
      "不向会员承诺股权、分红或平台估值增值",
      "金融/移民/医疗一律转介持牌方，不自营"],
     SAND)
card(s, Inches(8.95), Inches(1.5), Inches(3.9), Inches(5.2),
     "主办方能力画像",
     ["中东本地运营与合规经验",
      "华人高净值客户触达能力",
      "房产项目议价与交付管理",
      "跨语种内容与社群运营",
      "政商学资源协调（非替代经营）",
      "首年建议编制 20–24 人"],
     TEAL_LT)

# ============================================================================
# 6 背书建议（回答用户问题）
# ============================================================================
s = new("要不要大学 / 上海市政府 / 省级背书？", "结论：需要“辅助背书”，不需要“政府主办”", "背书")
rrect(s, Inches(0.45), Inches(1.5), Inches(12.4), Inches(1.15), TEAL)
txt(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(0.85),
    [("建议结论", 12, SAND, True, 4),
     ("大学与上海（或省级）相关机构的背书，对国内获客、出海论坛与品牌信任有明显加分；"
      "但迪拜落地的硬条件是阿联酋牌照、合规与履约能力。背书应定位为“指导 / 支持 / 学术合作”，"
      "避免形成“政府主办”预期，以免决策慢、承诺重、责任边界模糊。", 13, WHITE, False)])
rows = [
    ("大学研究机构", "推荐（阶段一）", "白皮书、论坛、择校/住房研究背书", "复旦住房研究中心、涉阿/国关院所、商学院出海研究团队"),
    ("上海市级平台", "可选（阶段二）", "上海侧出海活动、企业动员", "侨务/商务类平台、区级科创/企联、国际会客厅类载体"),
    ("省级层面", "谨慎（阶段三）", "大型峰会或区域合作时再谈", "避免早期过重行政叙事，优先商会与产业联盟"),
    ("驻外使领馆/商会", "强烈建议", "本地公信力与活动合规沟通", "驻迪拜总领馆相关活动、华人商会、行业协会"),
]
for i, (a, b, c, d) in enumerate(rows):
    y = Inches(2.85 + i * 0.95)
    rrect(s, Inches(0.45), y, Inches(12.4), Inches(0.85), INK_2, TEAL_LT)
    txt(s, Inches(0.65), y + Inches(0.18), Inches(2.4), Inches(0.5), [(a, 13, SAND_LT, True)], MSO_ANCHOR.MIDDLE)
    color = OK if "推荐" in b or "强烈" in b else (WARN if "可选" in b else MUTED)
    txt(s, Inches(3.2), y + Inches(0.18), Inches(2.2), Inches(0.5), [(b, 13, color, True)], MSO_ANCHOR.MIDDLE)
    txt(s, Inches(5.5), y + Inches(0.12), Inches(3.4), Inches(0.6), [(c, 12, OFF, False)], MSO_ANCHOR.MIDDLE)
    txt(s, Inches(9.0), y + Inches(0.12), Inches(3.6), Inches(0.6), [(d, 11, MUTED, False)], MSO_ANCHOR.MIDDLE)

# ============================================================================
# 7 三大战略
# ============================================================================
s = new("三大战略体系", "安居入口 → 创始会员 → 华人共同体", "战略")
pillars = [
    ("第一层 · 安居计划", SAND,
     ["帮助 1000 个华人家庭在迪拜安家",
      "精选 3–5 个战略合作项目",
      "精装 + 软装 + 智能家居 + 拎包入住",
      "交房后 45 天内可入住",
      "房产由 RERA 持牌方执行"]),
    ("第二层 · 创始会员", ACCENT,
     ["购房即成为创始会员（No.0001–1000）",
      "编号终身唯一、可继承、不可转售",
      "电子身份 + 创始金卡 + 权益包",
      "重大权益优先向创始会员开放",
      "企业会员单独编号序列 E001 起"]),
    ("第三层 · 华人共同体", TEAL_LT,
     ["1000 企业家 × 1000 家庭 × 1000 企业",
      "每月 ≥2 场线下活动（商务+家庭）",
      "每季度 1 次需求发布会",
      "12 个行业分会（会长任期 1 年）",
      "核心指标是成交与履约，不是人数"]),
]
for i, (title, accent, bullets) in enumerate(pillars):
    x = Inches(0.4 + i * 4.3)
    card(s, x, Inches(1.55), Inches(4.1), Inches(5.15), title, bullets, accent)

# ============================================================================
# 8 会员定位与价格（重点）
# ============================================================================
s = new("未来会员定位与价格体系", "规划假设以 AED 计价，立项时以实测替换", "定价")
draw_table(
    s,
    x0=Inches(0.4),
    y0=Inches(1.5),
    col_widths_in=[2.35, 2.15, 2.7, 1.2, 4.4],
    row_h_in=0.78,
    headers=["层级", "获取方式", "年费（规划）", "名额", "核心定位"],
    rows=[
        ["创始会员（个人）", "通过平台购房", "前3年免年费；第4年起 3,000", "1,000", "全权益 + 优先权 + 编号荣誉"],
        ["正式会员", "申请 + 审核", "AED 3,000 / 年", "不限", "标准生活权益，扩展层"],
        ["企业会员（标准）", "企业申请审核", "AED 12,000 / 年", "不限", "黄页 / 招聘 / 撮合 / 活动 2 席"],
        ["企业会员（战略）", "邀请制", "AED 60,000 / 年", "≤30", "联合品牌、峰会主论坛、深度撮合"],
        ["理事单位", "邀请 + 理事会", "AED 200,000 / 年", "≤12", "参与治理、优先合作、冠名权"],
    ],
    body_size=12,
    header_size=13,
)

# ============================================================================
# 9 会员分层画像
# ============================================================================
s = new("会员人群分层与获取", "A 核心买房家庭 → B 在地华人 → C 出海考察客群", "人群")
segs = [
    ("A 核心层", "1,000", "在迪拜购房/已购房华人家庭与企业主",
     ["核心诉求：安家、圈层、资产配置", "渠道：房产项目、渠道商、老带新", "产品：创始会员"]),
    ("B 扩展层", "10,000", "在阿联酋工作、创业的华人",
     ["核心诉求：生活服务、职业、合作", "渠道：社群、活动、内容", "产品：正式会员 / 企业会员"]),
    ("C 外围层", "100,000", "计划来中东考察/出海的中国企业与个人",
     ["核心诉求：信息、落地、对接", "渠道：线上内容、商会、展会", "产品：注册用户 → 转化入会"]),
]
for i, (name, n, who, bullets) in enumerate(segs):
    x = Inches(0.4 + i * 4.3)
    rrect(s, x, Inches(1.55), Inches(4.1), Inches(5.15), TEAL)
    txt(s, x + Inches(0.3), Inches(1.75), Inches(3.5), Inches(0.4),
        [(name, 16, SAND, True)])
    txt(s, x + Inches(0.3), Inches(2.25), Inches(3.5), Inches(0.55),
        [(n, 28, WHITE, True)])
    txt(s, x + Inches(0.3), Inches(2.9), Inches(3.5), Inches(0.7),
        [(who, 13, OFF, False)])
    txt(s, x + Inches(0.3), Inches(3.7), Inches(3.5), Inches(2.6),
        [("• " + b, 13, MUTED, False, 6) for b in bullets])

# ============================================================================
# 10 会员生态
# ============================================================================
s = new("针对未来会员的定位生态", "七大权益模块构成日常生活与商务底座", "生态")
rights = [
    ("安居", "物业/托管/维修/搬家/家政", "紧急维修 4h 响应"),
    ("医疗", "体检/绿通/家庭医生/保险", "持牌机构交付"),
    ("教育", "择校/中文/双语/升学", "年度择校白皮书"),
    ("餐饮生活", "餐厅/酒店/宴会折扣", "首年签约 100 家"),
    ("运动健康", "健身/球类/马拉松社群", "6 个常设运动社群"),
    ("商务出行", "用车/酒店/考察/签证", "每季商务考察团"),
    ("金融法律", "开户/财税/律所/注册", "每年 2 次免费初咨"),
]
for i, (a, b, c) in enumerate(rights):
    col, row = i % 4, i // 4
    if i == 6:
        x = Inches(0.45 + 1.5 * 3.2)
    else:
        x = Inches(0.45 + col * 3.2)
    y = Inches(1.55 + row * 2.55)
    rrect(s, x, y, Inches(3.05), Inches(2.3), TEAL)
    txt(s, x + Inches(0.2), y + Inches(0.35), Inches(2.65), Inches(1.7),
        [(a, 18, SAND_LT, True, 8),
         (b, 12, OFF, False, 8),
         (c, 12, SAND, False)])

# ============================================================================
# 11 六大中心
# ============================================================================
s = new("六大中心：生态供给侧", "首年只启 ①③⑥，避免全面铺开", "中心")
centers = [
    ("① 商业资源", "企业库 1000 / 撮合 100", "Y1", SAND),
    ("② 人才发展", "人才库 3000 / 推荐 100", "Y2", TEAL_LT),
    ("③ 品牌联盟", "品牌 50 / 集采 3", "Y1", SAND),
    ("④ 创业孵化", "项目 20 / 导师 30", "Y2", TEAL_LT),
    ("⑤ 投资发展", "路演 4 / 投资人 100", "Y3", MUTED),
    ("⑥ 公益文化", "公益 6 / 节日 2", "Y1", SAND),
]
for i, (name, kpi, phase, accent) in enumerate(centers):
    col, row = i % 3, i // 3
    x = Inches(0.45 + col * 4.25)
    y = Inches(1.55 + row * 2.55)
    rrect(s, x, y, Inches(4.05), Inches(2.3), TEAL)
    rect(s, x, y, Inches(0.12), Inches(2.3), accent)
    txt(s, x + Inches(0.35), y + Inches(0.35), Inches(3.4), Inches(1.6),
        [(name, 18, WHITE, True, 8),
         (kpi, 13, OFF, False, 8),
         (f"启动节奏：{phase}", 13, accent, True)])

# ============================================================================
# 12 社群搭建（重点）
# ============================================================================
s = new("社群搭建总览", "微信主阵地 + Facebook/海外矩阵协同", "社群")
card(s, Inches(0.4), Inches(1.5), Inches(4.15), Inches(5.2),
     "中国侧（微信生态）",
     ["公众号：政策/白皮书/会员故事，周 2 篇",
      "视频号/抖音：迪拜生活与避坑，周 4 条",
      "小红书：家庭决策（择校/亲子），周 5 条",
      "企业微信/私域：转化与留存，日更",
      "分层群：创始会员群 / 行业分会群 / 城市群",
      "规则：实名、禁广告、工单化客服"],
     SAND)
card(s, Inches(4.7), Inches(1.5), Inches(4.15), Inches(5.2),
     "海外侧（Facebook 等）",
     ["Facebook Group：迪拜华人公开层获客",
      "Facebook Page：品牌与活动传播",
      "Instagram：生活方式与现场感",
      "LinkedIn：企业会员与商务 B2B",
      "WhatsApp Community：本地即时触达",
      "小红书/抖音海外号：跨平台分发"],
     ACCENT)
card(s, Inches(9.0), Inches(1.5), Inches(3.9), Inches(5.2),
     "社群运营原则",
     ["不为热闹，为信任与成交",
      "每群设管理员 + 议题官",
      "活动后 48h 内发出名录与线索",
      "商务活动强制结构化连接",
      "双端账号统一会员 ID",
      "遵守 UAE PDPL + 中国 PIPL"],
     TEAL_LT)

# ============================================================================
# 13 社群矩阵详表
# ============================================================================
s = new("海外社群搭建清单", "渠道定位 · 内容 · 首年目标", "矩阵")
draw_table(
    s,
    x0=Inches(0.35),
    y0=Inches(1.42),
    col_widths_in=[2.0, 1.55, 3.05, 3.15, 2.85],
    row_h_in=0.52,
    headers=["渠道", "定位", "内容类型", "组织形态", "首年目标"],
    rows=[
        ["微信私域", "转化留存", "权益上新 / 答疑 / 活动", "创始群 1 + 分会 12", "续费率 ≥70%"],
        ["微信公众号", "深度权威", "政策解读 / 白皮书", "周更 2 篇", "粉丝 3 万"],
        ["Facebook 群", "海外获客", "本地生活 / 活动招募", "主群 1 + 城市分群", "成员 1 万"],
        ["Facebook 页", "品牌传播", "活动相册 / 媒体合作", "周更 3 条", "互动率 ≥3%"],
        ["Instagram", "生活方式", "现场 / 探店 / 会员故事", "周更 4 条", "粉丝 2 万"],
        ["LinkedIn", "企业商务", "峰会 / 合作案例 / 招聘", "周更 2 条", "企业线索 200"],
        ["WhatsApp", "本地触达", "活动提醒 / 紧急通知", "分层社区", "打开率 ≥60%"],
        ["线下分会", "关系沉淀", "月度沙龙 / 需求墙", "12 个行业分会", "月活出席 ≥80"],
    ],
    body_size=12,
    header_size=12,
)

# ============================================================================
# 14 年度活动
# ============================================================================
s = new("年度活动日历（旗舰节奏）", "每月有旗舰，每周有连接", "活动")
cal = [
    ("1月", "新年商业酒会"), ("2月", "春节庆典家庭日"), ("3月", "投资沙龙"),
    ("4月", "行业峰会①"), ("5月", "企业出海论坛"), ("6月", "半年度会员大会"),
    ("7–8月", "夏令营/线上维系"), ("9月", "中秋晚宴+慈善"), ("10月", "品牌/集采大会"),
    ("11月", "年度企业家论坛"), ("12月", "盛典+创始授勋"), ("常态", "分会月度沙龙"),
]
for i, (m, t) in enumerate(cal):
    col, row = i % 4, i // 4
    x = Inches(0.45 + col * 3.2)
    y = Inches(1.55 + row * 1.7)
    rrect(s, x, y, Inches(3.05), Inches(1.5), TEAL)
    txt(s, x + Inches(0.2), y + Inches(0.3), Inches(2.65), Inches(1.0),
        [(m, 14, SAND, True, 6), (t, 14, WHITE, False)])

# ============================================================================
# 15 商业模式
# ============================================================================
s = new("商业模式与收入演进", "第一年靠房产养平台，第三年去房产依赖", "商业")
card(s, Inches(0.4), Inches(1.5), Inches(6.2), Inches(5.2),
     "十大收入来源",
     ["房地产服务佣金 / 渠道分成",
      "房屋托管及物业增值",
      "会员年费（个人 / 企业 / 理事）",
      "品牌合作与赞助",
      "广告及推广",
      "商城与供应链（含集采）",
      "活动及会展（票务 / 展位 / 冠名）",
      "医疗教育旅游等服务分成",
      "金融保险法律转介分成",
      "企业出海及商务落地服务"],
     SAND)
# 收入结构表
txt(s, Inches(6.95), Inches(1.55), Inches(5.8), Inches(0.4),
    [("收入结构演进（规划）", 15, SAND_LT, True)])
struct = [
    ("类型", "Y1", "Y3", "Y5"),
    ("房地产相关", "75%", "45%", "30%"),
    ("会员及企业服务", "12%", "25%", "28%"),
    ("生态服务分成", "8%", "20%", "27%"),
    ("品牌/活动/广告", "5%", "10%", "15%"),
]
for i, row in enumerate(struct):
    y = Inches(2.15 + i * 0.7)
    for j, val in enumerate(row):
        x = Inches(6.95 + j * 1.45)
        bg = SAND if i == 0 else (TEAL if i % 2 else INK_2)
        fc = INK if i == 0 else WHITE
        rrect(s, x, y, Inches(1.38), Inches(0.58), bg)
        txt(s, x, y + Inches(0.12), Inches(1.38), Inches(0.4),
            [(val, 12, fc, i == 0 or j == 0)], align=PP_ALIGN.CENTER)
rrect(s, Inches(6.95), Inches(5.75), Inches(5.85), Inches(0.9), INK_2, ACCENT)
txt(s, Inches(7.15), Inches(5.9), Inches(5.5), Inches(0.65),
    [("战略含义：Y3 生态+会员收入合计应 ≥45%，否则平台只是房产中介另一张皮。",
      12, SAND_LT, False)])

# ============================================================================
# 16 财务
# ============================================================================
s = new("三年财务简化模型（假设）", "盈亏平衡约 143 套；启动资金建议 AED 5M", "财务")
fin = [
    ("项目", "第 1 年", "第 2 年", "第 3 年"),
    ("收入合计", "9.64M", "22.50M", "41.10M"),
    ("成本合计", "7.60M", "13.05M", "20.30M"),
    ("税前利润", "2.04M", "9.45M", "20.80M"),
]
for i, row in enumerate(fin):
    y = Inches(1.55 + i * 0.75)
    for j, val in enumerate(row):
        x = Inches(0.5 + j * 3.1)
        bg = SAND if i == 0 else (TEAL if i < 3 else INK_2)
        fc = INK if i == 0 else (OK if i == 3 else WHITE)
        rrect(s, x, y, Inches(2.95), Inches(0.65), bg)
        txt(s, x, y + Inches(0.15), Inches(2.95), Inches(0.4),
            [(val + (" AED" if i and j else ""), 14, fc, True)], align=PP_ALIGN.CENTER)
assumps = [
    ("首年成交", "200 套"),
    ("均价", "AED 1.8M"),
    ("佣金净得", "2%"),
    ("正式会员", "300 人"),
    ("企业会员", "50 家"),
    ("续费率", "70%→80%"),
]
for i, (a, b) in enumerate(assumps):
    col, row = i % 3, i // 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(4.8 + row * 0.95)
    rrect(s, x, y, Inches(4.0), Inches(0.82), TEAL)
    txt(s, x + Inches(0.25), y + Inches(0.2), Inches(3.5), Inches(0.45),
        [(f"{a}：{b}", 14, OFF, False)], anchor=MSO_ANCHOR.MIDDLE)

# ============================================================================
# 17 90天
# ============================================================================
s = new("首年启动：90 天冲刺", "奠基 → 搭建 → 引爆；设 Go / No-Go 闸门", "90天")
phases = [
    ("D1–30 奠基", ["法律架构与牌照方案", "2–3 个房产项目 LOI", "品牌 VI 与命名核准",
                  "核心团队 8 人 Offer", "财务模型与预算审批", "会员章程与权益手册"]),
    ("D31–60 搭建", ["公司注册与基础牌照", "签约首批 30 家合作商", "小程序 1 期启动",
                   "内容矩阵开播", "种子会员 50 位", "CRM / 工单上线"]),
    ("D61–90 引爆", ["发布会 200–300 人", "小程序上线可用", "创始会员正式招募",
                   "合作商扩至 80 家", "首场行业沙龙", "月度经营看板"]),
]
for i, (title, items) in enumerate(phases):
    x = Inches(0.4 + i * 4.3)
    card(s, x, Inches(1.5), Inches(4.1), Inches(4.0), title, items, SAND if i != 1 else ACCENT)
rrect(s, Inches(0.4), Inches(5.7), Inches(12.5), Inches(1.05), INK_2, OK)
txt(s, Inches(0.65), Inches(5.85), Inches(12.0), Inches(0.8),
    [("Go 标准：实体与房产合作落地 · 创始会员≥100 · 成交≥25 套 · 合作商≥80 · 小程序可用。"
      "任一严重不达标 → 暂停扩张，回到验证阶段。", 13, OFF, False)])

# ============================================================================
# 18 KPI
# ============================================================================
s = new("KPI 与北极星指标", "MAM × ARPU；用数据驱动续费与生态", "KPI")
kpis = [
    ("增长", ["创始会员 1000（18 月内）", "房产成交 200 套", "企业会员 50 家"]),
    ("留存", ["续费率 ≥70%", "NPS ≥50", "月活跃率 ≥40%"]),
    ("生态", ["认证合作商 150", "人均权益模块 ≥3", "撮合成功 100"]),
    ("运营/财务", ["活动 40 场", "投诉 48h 解决 ≥95%", "LTV/CAC ≥3"]),
]
for i, (t, lines) in enumerate(kpis):
    x = Inches(0.4 + i * 3.25)
    card(s, x, Inches(1.55), Inches(3.1), Inches(5.1), t, lines, SAND if i % 2 == 0 else ACCENT)

# ============================================================================
# 19 合规
# ============================================================================
s = new("合规红线与关键风险", "先合法，再扩张", "合规")
card(s, Inches(0.4), Inches(1.5), Inches(6.2), Inches(5.2),
     "合规铁律",
     ["不承诺投资回报 / 租金保证（除非书面担保并披露）",
      "不吸收资金、不代持、不设资金池",
      "不做未持牌移民 / 法律 / 医疗 / 金融经营",
      "不用未经证实的“最大/第一/保证”宣传",
      "会员数据未经授权不得对外提供",
      "活动尊重当地宗教文化与斋月安排",
      "遵守 UAE PDPL 与中国 PIPL 跨境规则"],
     ACCENT)
card(s, Inches(6.85), Inches(1.5), Inches(6.0), Inches(5.2),
     "高优先级风险",
     ["R1 房产销售不及预期 → 现金流断裂",
      "R2 过度依赖房产收入 → Y2 强制生态≥35%",
      "R3 权益承诺无法兑付 → 先签合作商再上线",
      "R4 牌照/合规问题 → 立项前法律意见书",
      "R5 创始身份贬值 → 治理权 + 透明度报告",
      "R8 品牌舆情 → 24h 响应与发言人制度",
      "盈亏平衡点约 143 套，必须月度监控"],
     SAND)

# ============================================================================
# 20 落地清单（融合四重点）
# ============================================================================
s = new("落地执行抓手（四重点融合）", "把战略拆成可分配的工作包", "落地")
packs = [
    ("主办方", ["注册双主体与牌照", "理事会与章程", "团队 20–24 人编制", "预算与启动金 5M AED"]),
    ("会员定价", ["五级产品上架", "权益手册 + SLA", "入会 7 天顾问接待", "续费与升级路径"]),
    ("会员生态", ["首年启 ①③⑥ 中心", "合作商准入五步", "首年 150 家认证商", "保障金 3% 计提"]),
    ("社群搭建", ["微信私域分层", "FB 主群+主页", "IG/LinkedIn/WhatsApp", "12 分会与活动日历"]),
]
for i, (t, lines) in enumerate(packs):
    x = Inches(0.4 + i * 3.25)
    card(s, x, Inches(1.55), Inches(3.1), Inches(5.1), t, lines,
         [SAND, ACCENT, TEAL_LT, OK][i])

# ============================================================================
# 21 背书行动建议
# ============================================================================
s = new("背书辅助落地：分阶段动作", "先商会与学术，后市区级平台，省级慎用", "行动")
steps = [
    ("阶段一 · 0–90 天",
     ["对接驻迪拜总领馆相关活动渠道与华人商会",
      "签约 1–2 家大学/研究机构为学术支持单位",
      "上海侧以区级企联/国际会客厅做活动共创",
      "对外话术：联合支持 / 学术指导，非政府主办"]),
    ("阶段二 · 3–12 月",
     ["白皮书联名发布，抬升国内信任度",
      "争取市级相关平台作为“支持单位”出现在峰会",
      "形成可复制的上海—迪拜双向考察团产品",
      "评估是否引入省级产业联盟作战略合作"]),
    ("阶段三 · 12 月后",
     ["在有实质 GMV 与会员履约数据后再谈更高级背书",
      "背书服务于扩张，不替代产品与合规",
      "保持商业主体独立决策与快速迭代能力",
      "所有背书关系写入合作备忘录，明确权责"]),
]
for i, (t, lines) in enumerate(steps):
    x = Inches(0.4 + i * 4.3)
    card(s, x, Inches(1.55), Inches(4.1), Inches(5.15), t, lines, SAND if i == 0 else (ACCENT if i == 1 else TEAL_LT))

# ============================================================================
# 22 结束页
# ============================================================================
PAGE += 1
s = slide()
rect(s, 0, 0, SW, Pt(5), SAND)
rect(s, 0, Inches(5.9), SW, Pt(2), SAND)
txt(s, Inches(1.0), Inches(1.8), Inches(11.3), Inches(2.4),
    [("一套房，只能改变一个家庭的居住方式；", 20, OFF, False, 10),
     ("一个平台，却可以改变一代华人在中东的发展方式。", 20, WHITE, True, 16),
     ("中东华人共创共享联盟", 28, SAND, True, 8),
     ("让世界因连接而更有价值，让华人因共创而更加精彩。", 16, SAND_LT, False)])
txt(s, Inches(1.0), Inches(5.0), Inches(11.3), Inches(0.7),
    [("下一步：法律意见书 → 真实财务参数替换 → 立项评审确认 90 天责任人与预算",
      13, MUTED, False)])
footer(s, PAGE)

# ============================================================================
# 保存
# ============================================================================
OUT_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "deliverables")
os.makedirs(OUT_DIR, exist_ok=True)
OUT = os.path.join(OUT_DIR, "中东华人共创共享联盟_战略与落地执行.pptx")
prs.save(OUT)
print(f"已生成: {OUT}  共 {len(prs.slides)} 页")
