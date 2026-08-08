#!/usr/bin/env python3
"""生成《接手失败项目的通用框架》方法论 PPT。"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# 视觉方向：墨色骨架 + 青绿行动色（避免紫系/奶油陶土系）
INK = RGBColor(0x0F, 0x17, 0x2A)
INK_SOFT = RGBColor(0x1E, 0x29, 0x3B)
TEAL = RGBColor(0x0F, 0x76, 0x6E)
TEAL_LIGHT = RGBColor(0x14, 0xB8, 0xA6)
AMBER = RGBColor(0xD9, 0x77, 0x06)
SLATE = RGBColor(0x47, 0x55, 0x69)
SOFT = RGBColor(0xF1, 0xF5, 0xF9)
MINT = RGBColor(0xEC, 0xFD, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x29, 0x3B)
LINE = RGBColor(0xCB, 0xD5, 0xE1)

FONT = "Microsoft YaHei"
FOOTER = "接手失败项目的通用框架 · 全量摸底 → 按己逻辑重构 → 完成后及时转向"
TOTAL = 12


def add_rect(slide, x, y, w, h, fill, *, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_round(slide, x, y, w, h, fill, *, adj=0.08):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    shp.adjustments[0] = adj
    return shp


def set_run(run, text, *, size=16, bold=False, color=DARK, font=FONT):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    size=16,
    bold=False,
    color=DARK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        set_run(p.add_run(), line, size=size, bold=bold, color=color)
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=14, color=DARK, mark="●", mark_color=TEAL, spacing=1.2):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Emu(0))
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = spacing
        set_run(p.add_run(), f"{mark}  ", size=size, color=mark_color)
        set_run(p.add_run(), item, size=size, color=color)
    return tb


def header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SW, Inches(0.9), INK)
    add_rect(slide, 0, Inches(0.9), SW, Inches(0.06), TEAL_LIGHT)
    add_text(slide, Inches(0.5), Inches(0.18), Inches(12.2), Inches(0.4), title, size=24, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.52), Inches(12.2), Inches(0.3), subtitle, size=12, color=TEAL_LIGHT)


def footer(slide, page):
    add_text(slide, Inches(0.5), Inches(7.1), Inches(10.8), Inches(0.28), FOOTER, size=10, color=SLATE)
    add_text(
        slide,
        Inches(11.4),
        Inches(7.1),
        Inches(1.4),
        Inches(0.28),
        f"{page} / {TOTAL}",
        size=10,
        color=SLATE,
        align=PP_ALIGN.RIGHT,
    )


def card(slide, x, y, w, h, title, items, *, accent=TEAL, title_color=None):
    add_round(slide, x, y, w, h, SOFT)
    add_rect(slide, x, y, Inches(0.1), h, accent)
    add_text(
        slide,
        x + Inches(0.28),
        y + Inches(0.15),
        w - Inches(0.4),
        Inches(0.35),
        title,
        size=15,
        bold=True,
        color=title_color or accent,
    )
    add_bullets(
        slide,
        x + Inches(0.28),
        y + Inches(0.55),
        w - Inches(0.45),
        h - Inches(0.7),
        items,
        size=12,
        mark_color=accent,
    )


def build_ppt(output_path: Path) -> None:
    global SW, SH
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    # 1 封面
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, INK)
    add_rect(s, 0, 0, Inches(0.2), SH, TEAL_LIGHT)
    add_text(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.35), "AI 项目管理 · 产品接手 · 可复用方法论", size=14, color=TEAL_LIGHT)
    add_text(
        s,
        Inches(0.8),
        Inches(2.0),
        Inches(11.5),
        Inches(1.4),
        "接手失败项目的\n通用框架",
        size=42,
        bold=True,
        color=WHITE,
    )
    add_rect(s, Inches(0.8), Inches(3.6), Inches(2.4), Inches(0.06), TEAL_LIGHT)
    add_text(
        s,
        Inches(0.8),
        Inches(3.9),
        Inches(11.5),
        Inches(0.9),
        "先全量摸底 → 再按自己的逻辑重构 → 完成后及时转向\n把一次「救火」沉淀成可迁移的作战手册",
        size=16,
        color=SOFT,
    )
    steps = [("01", "全量摸底"), ("02", "逻辑重构"), ("03", "及时转向")]
    for i, (num, name) in enumerate(steps):
        x = Inches(0.8 + i * 3.5)
        add_round(s, x, Inches(5.3), Inches(3.2), Inches(1.0), INK_SOFT)
        add_text(s, x + Inches(0.2), Inches(5.45), Inches(1.0), Inches(0.35), num, size=18, bold=True, color=TEAL_LIGHT)
        add_text(s, x + Inches(1.1), Inches(5.55), Inches(1.9), Inches(0.4), name, size=18, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    # 2 为什么需要框架
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "为什么「接手失败项目」需要框架", "多数失败不是能力不够，而是缺一套可重复的接手节奏")
    pains = [
        ("盲目救火", "一上来就改功能、修 Bug，越修越散，无法建立全局图景"),
        ("继承幻觉", "默认沿用前任的目标、架构与叙事，把错误路径继续走下去"),
        ("无限收尾", "清理没有终点，接手变成长期陪葬，错过下一阶段价值创造"),
        ("经验不可迁移", "每次都靠个人直觉，组织层面无法沉淀为可复用能力"),
    ]
    for i, (t, b) in enumerate(pains):
        y = Inches(1.25 + i * 1.3)
        add_round(s, Inches(0.5), y, Inches(12.3), Inches(1.15), SOFT if i % 2 == 0 else MINT)
        add_text(s, Inches(0.8), y + Inches(0.2), Inches(2.5), Inches(0.35), t, size=16, bold=True, color=TEAL)
        add_text(s, Inches(3.5), y + Inches(0.25), Inches(8.8), Inches(0.65), b, size=14, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 2)

    # 3 总览
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "框架总览：三步闭环", "摸底建立真相 → 重构建立主权 → 转向释放产能")
    phases = [
        (TEAL, "01 全量摸底", "建立事实底座", ["资产 / 债务 / 干系人", "可用 / 不可用 / 未知", "不做局部修补"]),
        (INK, "02 逻辑重构", "用自己的操作系统接管", ["目标重定义", "杀 / 留 / 重写", "叙事与决策权归己"]),
        (AMBER, "03 及时转向", "定义完成并离开清理态", ["写清「接手完成」标准", "交接或进入创造期", "沉淀可复用 playbook"]),
    ]
    for i, (accent, title, sub, items) in enumerate(phases):
        x = Inches(0.45 + i * 4.25)
        add_round(s, x, Inches(1.3), Inches(4.0), Inches(5.2), SOFT)
        add_rect(s, x, Inches(1.3), Inches(4.0), Inches(1.15), accent)
        add_text(s, x + Inches(0.25), Inches(1.4), Inches(3.5), Inches(0.45), title, size=18, bold=True, color=WHITE)
        add_text(s, x + Inches(0.25), Inches(1.85), Inches(3.5), Inches(0.35), sub, size=12, color=SOFT)
        add_bullets(s, x + Inches(0.3), Inches(2.7), Inches(3.4), Inches(3.3), items, size=14, mark_color=accent)
        if i < 2:
            add_text(s, x + Inches(3.7), Inches(3.5), Inches(0.5), Inches(0.4), "→", size=22, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
    footer(s, 3)

    # 4 步骤一详解
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "步骤一｜全量摸底", "先看见全貌，再谈动作；摸底不是拖延，是防踩雷")
    card(
        s,
        Inches(0.4),
        Inches(1.25),
        Inches(4.0),
        Inches(5.3),
        "摸什么",
        [
            "目标与承诺：当初卖什么、现在还剩什么",
            "资产清单：代码、数据、文档、账号、合同",
            "债务清单：技术债、合规债、信誉债",
            "干系人地图：谁能拍板、谁在等结果",
            "外部约束：时间、预算、法务、品牌",
        ],
        accent=TEAL,
    )
    card(
        s,
        Inches(4.65),
        Inches(1.25),
        Inches(4.0),
        Inches(5.3),
        "怎么摸",
        [
            "只读优先：先读后改，禁止边摸边大改",
            "三色标注：绿可用 / 黄存疑 / 红不可用",
            "事实与观点分离：记录证据，不下结论过早",
            "时间盒：例如 3–5 天完成第一轮全量盘点",
            "输出物：一页现状真相图 + 风险清单",
        ],
        accent=INK,
    )
    card(
        s,
        Inches(8.9),
        Inches(1.25),
        Inches(4.0),
        Inches(5.3),
        "摸底完成标准",
        [
            "能讲清「为什么失败」的主因假设",
            "能指出可 salvage 的 20% 高价值资产",
            "能列出必须立刻停掉的危险动作",
            "干系人对「现状陈述」基本无异议",
            "进入重构前，团队知道边界在哪",
        ],
        accent=AMBER,
    )
    footer(s, 4)

    # 5 摸底清单示意
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "摸底输出：一页「现状真相图」", "把复杂项目压成可决策的四象限")
    quads = [
        (Inches(0.45), Inches(1.25), TEAL, "可直接复用", ["稳定核心模块", "已验证的数据/流程", "可迁移的客户关系"]),
        (Inches(6.85), Inches(1.25), AMBER, "需改造后使用", ["半成品功能", "文档不全但有价值", "依赖脆弱的集成"]),
        (Inches(0.45), Inches(4.15), INK, "应立即冻结/停用", ["无人维护的旁路", "合规风险点", "虚假进度与空壳指标"]),
        (Inches(6.85), Inches(4.15), SLATE, "信息不足 · 待查", ["关键决策无人能答", "历史承诺口径不一", "成本/收益无数据"]),
    ]
    for x, y, c, title, items in quads:
        add_round(s, x, y, Inches(6.0), Inches(2.65), SOFT)
        add_rect(s, x, y, Inches(6.0), Inches(0.5), c)
        add_text(s, x + Inches(0.25), y + Inches(0.08), Inches(5.5), Inches(0.35), title, size=16, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_bullets(s, x + Inches(0.3), y + Inches(0.7), Inches(5.4), Inches(1.7), items, size=13, mark_color=c)
    footer(s, 5)

    # 6 步骤二详解
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "步骤二｜按自己的逻辑重构", "接手不是继承前任世界观，而是重建可执行的操作系统")
    principles = [
        ("重定义成功", "用当前约束重写目标：什么叫「救回来」、什么叫「可以放弃」"),
        ("先杀后留", "默认不信任旧路线；能删则删，能绕则绕，重写只投在主路径"),
        ("主权在己", "决策口径、优先级、验收标准由接手方定义，避免「双轨指挥」"),
        ("小闭环验证", "用最短闭环证明新逻辑成立，再扩大重构半径"),
        ("叙事同步改", "对内对外统一新故事：失败归因、当前策略、下一步里程碑"),
        ("工具服从逻辑", "AI / 流程 / 看板只是执行层，先定逻辑再选工具"),
    ]
    for i, (t, b) in enumerate(principles):
        col, row = i % 3, i // 3
        x = Inches(0.4 + col * 4.25)
        y = Inches(1.25 + row * 2.7)
        add_round(s, x, y, Inches(4.05), Inches(2.4), MINT if row == 0 else SOFT)
        add_text(s, x + Inches(0.25), y + Inches(0.3), Inches(3.5), Inches(0.4), t, size=16, bold=True, color=TEAL)
        add_text(s, x + Inches(0.25), y + Inches(0.9), Inches(3.5), Inches(1.2), b, size=13, color=DARK)
    footer(s, 6)

    # 7 杀留重写
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "重构决策矩阵：杀 / 留 / 重写", "把情绪争论变成可执行的分类动作")
    headers_row = ["动作", "适用信号", "典型动作", "风险若做错"]
    rows = [
        ["杀", "无主责、无证据、无用户", "下线、冻结、归档、切断入口", "继续烧资源维持空壳"],
        ["留", "稳定、可证伪、低维护成本", "最小改动纳入新边界，写清契约", "误删核心资产"],
        ["重写", "主路径但旧设计不可救", "按新逻辑重做最小可行切片", "大爆炸重写失控"],
    ]
    # manual table-like cards
    add_round(s, Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.55), INK)
    for i, h in enumerate(headers_row):
        add_text(s, Inches(0.55 + i * 3.1), Inches(1.28), Inches(3.0), Inches(0.4), h, size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    colors = [AMBER, TEAL, INK]
    for r, row in enumerate(rows):
        y = Inches(1.9 + r * 1.5)
        add_round(s, Inches(0.4), y, Inches(12.5), Inches(1.35), SOFT)
        add_rect(s, Inches(0.4), y, Inches(0.12), Inches(1.35), colors[r])
        for i, cell in enumerate(row):
            bold = i == 0
            add_text(
                s,
                Inches(0.55 + i * 3.1),
                y + Inches(0.4),
                Inches(3.0),
                Inches(0.55),
                cell,
                size=14,
                bold=bold,
                color=colors[r] if i == 0 else DARK,
                anchor=MSO_ANCHOR.MIDDLE,
            )
    footer(s, 7)

    # 8 步骤三详解
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "步骤三｜完成后及时转向", "接手有终点；清理态不是事业态")
    add_round(s, Inches(0.4), Inches(1.25), Inches(12.5), Inches(1.5), MINT)
    add_text(
        s,
        Inches(0.7),
        Inches(1.5),
        Inches(12.0),
        Inches(1.0),
        "核心原则：提前写下「接手完成」的可验收定义。没有完成定义，就会永远停在修补循环里。",
        size=16,
        bold=True,
        color=TEAL,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    exits = [
        ("完成定义示例", ["主路径可演示且可复现", "关键风险已封堵或有 owner", "干系人接受新目标与边界", "文档与职责完成交接"]),
        ("转向去向", ["进入增长/交付创造期", "交给稳态 Owner 运营", "正式关停并沉淀复盘", "拆成新项目重新立项"]),
        ("防拖延机制", ["公开倒计时与里程碑", "每周问：还在清理还是在创造？", "清理任务设 WIP 上限", "到期未完成则升级决策"]),
    ]
    for i, (t, items) in enumerate(exits):
        x = Inches(0.4 + i * 4.25)
        card(s, x, Inches(3.0), Inches(4.05), Inches(3.5), t, items, accent=TEAL if i != 2 else AMBER)
    footer(s, 8)

    # 9 反模式
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "常见反模式（务必避开）", "框架的价值一半在「做什么」，一半在「坚决不做什么」")
    antis = [
        ("边摸边大改", "事实未稳就重构，导致无法归因：是旧问题还是新动作造成的？"),
        ("讨好式继承", "不敢否定前任目标，结果变成「双目标并行」，资源被撕碎"),
        ("完美主义清理", "追求把所有债还清再前进，错过窗口期"),
        ("工具崇拜", "先上 AI Agent / 看板 / 新框架，却没有接手逻辑"),
        ("口头交接", "没有书面完成定义与 owner，转向后问题回流到接手人"),
        ("英雄主义", "个人救火替代系统建设，经验无法复制，人一走项目再崩"),
    ]
    for i, (t, b) in enumerate(antis):
        col, row = i % 3, i // 3
        x = Inches(0.4 + col * 4.25)
        y = Inches(1.25 + row * 2.7)
        add_round(s, x, y, Inches(4.05), Inches(2.4), SOFT)
        add_rect(s, x, y, Inches(4.05), Inches(0.55), AMBER if row == 0 else INK)
        add_text(s, x + Inches(0.2), y + Inches(0.1), Inches(3.6), Inches(0.35), t, size=15, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.25), y + Inches(0.8), Inches(3.55), Inches(1.35), b, size=13, color=DARK)
    footer(s, 9)

    # 10 AI 项目管理场景
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "场景 A｜AI 项目管理", "用框架接管「跑偏的 AI 项目」而不是继续堆 Prompt")
    mapping = [
        ("全量摸底", ["盘点：目标指标、数据质量、模型/Agent 边界、评测集、成本账单", "分清：演示效果 vs 可上线能力", "标出幻觉高发路径与无人负责模块"]),
        ("逻辑重构", ["重定：业务问题 → 任务分解 → 人机分工", "杀无效 Agent 链路；留可评测链路；重写主任务编排", "建立：验收用例、失败回退、人工兜底"]),
        ("及时转向", ["主链路达标即冻结探索性实验", "把探索预算单列，避免污染交付", "沉淀：提示词/评测/运维 playbook 后交稳态团队"]),
    ]
    for i, (t, items) in enumerate(mapping):
        x = Inches(0.4 + i * 4.25)
        add_round(s, x, Inches(1.25), Inches(4.05), Inches(5.3), SOFT)
        add_rect(s, x, Inches(1.25), Inches(4.05), Inches(0.7), TEAL)
        add_text(s, x + Inches(0.2), Inches(1.38), Inches(3.6), Inches(0.45), t, size=18, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_bullets(s, x + Inches(0.25), Inches(2.2), Inches(3.55), Inches(4.0), items, size=13)
    footer(s, 10)

    # 11 产品接手场景
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "场景 B｜产品接手", "新产品负责人如何在 2–4 周内夺回方向权")
    weeks = [
        ("第 1 周 · 摸底", ["用户真实路径走查", "收入/留存/成本三张表", "干系人一对一访谈", "输出真相图与停做清单"]),
        ("第 2–3 周 · 重构", ["重写产品北极星与非目标", "砍路线图到一条主路径", "重排团队职责与节奏", "对外统一新叙事"]),
        ("第 4 周 · 转向", ["演示可验收切片", "明确 Owner 与 SLA", "关闭接手专项看板", "进入常态迭代/增长"]),
    ]
    for i, (t, items) in enumerate(weeks):
        x = Inches(0.4 + i * 4.25)
        card(s, x, Inches(1.25), Inches(4.05), Inches(5.3), t, items, accent=TEAL if i != 2 else AMBER)
    footer(s, 11)

    # 12 一页作战手册
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, INK)
    add_rect(s, 0, 0, Inches(0.2), SH, TEAL_LIGHT)
    add_text(s, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.5), "一页作战手册", size=28, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(1.05), Inches(11.5), Inches(0.4), "接到失败项目时，按此顺序行动，不要跳步。", size=14, color=TEAL_LIGHT)
    lines = [
        "1. 设时间盒，宣布「只读摸底期」——禁止大改。",
        "2. 产出真相图：可复用 / 需改造 / 冻结 / 待查。",
        "3. 重写成功定义与非目标，获关键干系人书面确认。",
        "4. 用杀/留/重写矩阵压缩范围到一条主路径。",
        "5. 跑通最短可验证闭环，再扩大重构半径。",
        "6. 写下接手完成标准；到期转向，清理任务设上限。",
        "7. 沉淀 playbook，交给稳态 Owner，自己进入下一创造。",
    ]
    for i, line in enumerate(lines):
        add_text(s, Inches(0.8), Inches(1.6 + i * 0.55), Inches(11.5), Inches(0.45), line, size=16, color=SOFT)
    add_text(
        s,
        Inches(0.8),
        Inches(6.5),
        Inches(11.5),
        Inches(0.5),
        "适用：AI 项目管理 · 产品接手 · 技术债抢救 · 失败业务盘活",
        size=13,
        color=TEAL_LIGHT,
    )

    prs.save(str(output_path))
    print(f"已生成: {output_path}")


if __name__ == "__main__":
    out = Path(__file__).resolve().parents[1] / "deliverables" / "接手失败项目的通用框架.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    build_ppt(out)
