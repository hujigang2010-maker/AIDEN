#!/usr/bin/env python3
"""生成《泰隆银行上海分行与腾讯 WorkBuddy 合作推进情况汇报》精简 PPT（5页）。"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt


# 银行汇报风格：深蓝主色 + 金色点缀
NAVY = RGBColor(0x0B, 0x2F, 0x5B)
NAVY_DARK = RGBColor(0x06, 0x1E, 0x3C)
GOLD = RGBColor(0xC4, 0xA3, 0x5A)
LIGHT = RGBColor(0xF3, 0xF6, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1F, 0x2A, 0x37)
GRAY = RGBColor(0x5C, 0x6B, 0x7A)
SOFT = RGBColor(0xE8, 0xEE, 0xF5)
ACCENT_RED = RGBColor(0xA6, 0x3D, 0x2F)

FONT = "Microsoft YaHei"
FOOTER_TEXT = "泰隆银行上海分行 × 腾讯 WorkBuddy · 合作推进情况汇报 · 2026-08-03"


def add_rect(slide, x, y, w, h, fill, *, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_round(slide, x, y, w, h, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    shp.adjustments[0] = 0.08
    return shp


def set_run(run, text, *, size=16, bold=False, color=DARK, font=FONT):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    size=16,
    bold=False,
    color=DARK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font=FONT,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        set_run(p.add_run(), line, size=size, bold=bold, color=color, font=font)
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=15, color=DARK, bullet_color=GOLD, spacing=1.25):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = spacing
        set_run(p.add_run(), "●  ", size=size, color=bullet_color)
        set_run(p.add_run(), item, size=size, color=color)
    return tb


def header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SW, Inches(0.95), NAVY)
    add_rect(slide, 0, Inches(0.95), SW, Inches(0.05), GOLD)
    add_text(slide, Inches(0.5), Inches(0.2), Inches(12.2), Inches(0.45), title, size=24, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.58), Inches(12.2), Inches(0.3), subtitle, size=12, color=LIGHT)


def footer(slide, page, total=5):
    add_text(slide, Inches(0.5), Inches(7.12), Inches(10.5), Inches(0.28), FOOTER_TEXT, size=10, color=GRAY)
    add_text(
        slide,
        Inches(11.5),
        Inches(7.12),
        Inches(1.3),
        Inches(0.28),
        f"{page} / {total}",
        size=10,
        color=GRAY,
        align=PP_ALIGN.RIGHT,
    )


def build_ppt(output_path: Path) -> None:
    global SW, SH
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    # ========== 第1页：封面 / 总体判断 ==========
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, NAVY_DARK)
    add_rect(s, 0, 0, Inches(0.18), SH, GOLD)
    add_text(s, Inches(0.8), Inches(1.35), Inches(11.5), Inches(0.4), "行长专题汇报 · 内部材料", size=14, color=GOLD)
    add_text(
        s,
        Inches(0.8),
        Inches(1.9),
        Inches(11.5),
        Inches(1.3),
        "泰隆银行上海分行与腾讯 WorkBuddy\n合作推进情况汇报",
        size=36,
        bold=True,
        color=WHITE,
    )
    add_rect(s, Inches(0.8), Inches(3.45), Inches(2.2), Inches(0.05), GOLD)
    add_text(
        s,
        Inches(0.8),
        Inches(3.75),
        Inches(11.5),
        Inches(1.3),
        "主方向：积分商城兑换 + 开户赠送AI权益；同步补充「行员采购 + 行外采购」双轨份额，\n"
        "并研究联名信用卡。建议原则同意推进，1,500万元仅为方案框架，先合规审查与小范围试点。",
        size=15,
        color=LIGHT,
    )
    add_text(
        s,
        Inches(0.8),
        Inches(5.5),
        Inches(11.5),
        Inches(0.4),
        "汇报对象：泰隆银行上海分行行长　　汇报日期：2026年8月3日",
        size=14,
        color=GOLD,
    )
    add_text(
        s,
        Inches(0.8),
        Inches(6.2),
        Inches(11.5),
        Inches(0.35),
        "三条主线：行外客户｜行内员工｜联名信用卡　　状态：已达成共识｜初步方案｜尚待确认",
        size=13,
        color=GRAY,
    )

    # ========== 第2页：合作方向与模式 ==========
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "一、合作方向与拟议模式", "行外客户权益 + 行内员工份额 + 联名信用卡")

    # 三条主线标题行
    tracks = [
        (NAVY, GOLD, "行外客户（已达成共识）", "积分商城兑换 · 开户/拜访赠礼 · 联合推广"),
        (SOFT, NAVY, "行内员工（初步方案）", "办公提效系统/小程序 · 员工AI福利发放"),
        (SOFT, NAVY, "联名信用卡（尚待确认）", "泰隆银行 × WorkBuddy 联名卡获客"),
    ]
    for i, (bg, title_c, title, sub) in enumerate(tracks):
        x = Inches(0.4 + i * 4.25)
        add_round(s, x, Inches(1.2), Inches(4.05), Inches(1.15), bg)
        add_text(s, x + Inches(0.2), Inches(1.35), Inches(3.65), Inches(0.4), title, size=14, bold=True, color=title_c if bg != NAVY else GOLD)
        add_text(
            s,
            x + Inches(0.2),
            Inches(1.8),
            Inches(3.65),
            Inches(0.4),
            sub,
            size=11,
            color=WHITE if bg == NAVY else DARK,
        )

    cards = [
        ("行外", "积分商城 / 开户赠礼", "集中采购上架兑换；开户拜访以AI权益替代部分实物礼品；兑换码自助激活"),
        ("行内①", "办公提效系统/小程序", "与WorkBuddy联手打造内部办公提效工具，覆盖起草、纪要、知识检索等场景"),
        ("行内②", "员工AI福利发放", "以福利形式向行员发放账号/额度，提升个人AI使用能力；单独预留员工份额"),
        ("对外", "联名信用卡", "原则同意研究联名卡方案，AI权益绑定办卡激活与消费，扩大零售获客"),
    ]
    for i, (tag, title, body) in enumerate(cards):
        x = Inches(0.4 + (i % 4) * 3.2)
        y = Inches(2.6)
        dark = i == 0
        add_round(s, x, y, Inches(3.05), Inches(4.0), NAVY if dark else SOFT)
        add_text(
            s,
            x + Inches(0.18),
            y + Inches(0.25),
            Inches(2.7),
            Inches(0.35),
            tag,
            size=12,
            bold=True,
            color=GOLD if dark else NAVY,
        )
        add_text(
            s,
            x + Inches(0.18),
            y + Inches(0.7),
            Inches(2.7),
            Inches(0.9),
            title,
            size=18,
            bold=True,
            color=WHITE if dark else NAVY,
        )
        add_text(
            s,
            x + Inches(0.18),
            y + Inches(1.8),
            Inches(2.7),
            Inches(1.9),
            body,
            size=13,
            color=LIGHT if dark else DARK,
        )
    footer(s, 2)

    # ========== 第3页：采购框架与预期价值 ==========
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "二、采购方案设想与预期价值", "行员采购 + 行外采购双轨；1,500万为测算框架，非采购承诺")

    # 左：采购框架
    add_round(s, Inches(0.45), Inches(1.25), Inches(6.3), Inches(5.5), SOFT)
    add_text(s, Inches(0.7), Inches(1.45), Inches(5.8), Inches(0.4), "采购框架（初步方案）", size=18, bold=True, color=NAVY)
    metrics = [
        ("最高约1,500万", "总体方案设计上限；三阶段×约500万"),
        ("首期300–500万", "明确行外/行员份额比例后再审批"),
        ("双轨拆分核算", "行外：客户权益；行员：提效+福利"),
    ]
    for i, (big, small) in enumerate(metrics):
        y = Inches(2.0 + i * 1.0)
        add_round(s, Inches(0.75), y, Inches(5.7), Inches(0.88), WHITE)
        add_rect(s, Inches(0.75), y, Inches(0.12), Inches(0.88), GOLD)
        add_text(s, Inches(1.1), y + Inches(0.1), Inches(5.1), Inches(0.35), big, size=17, bold=True, color=NAVY)
        add_text(s, Inches(1.1), y + Inches(0.48), Inches(5.1), Inches(0.3), small, size=12, color=GRAY)

    add_text(
        s,
        Inches(0.75),
        Inches(5.2),
        Inches(5.7),
        Inches(1.3),
        "尚待腾讯提供：行外/行员可分配数量、有效期、阶梯折扣、\n未激活处理、付款验收；另需内部系统开发范围及联名卡权益包。",
        size=12,
        color=DARK,
    )

    # 右：预期价值
    add_round(s, Inches(7.0), Inches(1.25), Inches(5.85), Inches(5.5), NAVY)
    add_text(s, Inches(7.3), Inches(1.45), Inches(5.3), Inches(0.4), "预期价值", size=18, bold=True, color=GOLD)
    values = [
        "带动新增对公客户：差异化开户/拜访权益",
        "提升存量客户黏性与积分活跃",
        "带动结算与资金沉淀：商城/续费走我行账户",
        "提升行内办公效率与员工AI能力",
        "联名信用卡获客与品牌曝光",
        "塑造“金融+经营赋能”科技服务品牌",
    ]
    add_bullets(s, Inches(7.3), Inches(2.05), Inches(5.3), Inches(4.3), values, size=13, color=WHITE, bullet_color=GOLD, spacing=1.45)
    footer(s, 3)

    # ========== 第4页：风险与共识 ==========
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "三、风险事项与共识分层", "外部客户优先；行内提效与联名卡须专项论证")

    # 左侧风险
    add_text(s, Inches(0.5), Inches(1.2), Inches(6.2), Inches(0.35), "尚待确认 · 主要风险", size=16, bold=True, color=ACCENT_RED)
    risks = [
        ("安全合规", "行外SaaS需审查数据隔离与权限；行内系统/小程序合规要求更高，须专项评估"),
        ("采购成本", "缺兑换/激活/使用数据，忌一次性大额采购；行外与行员份额分账核算"),
        ("技术交付", "兑换码自助激活、内部系统对接、员工发放流程均需完整测试"),
        ("联名信用卡", "卡种准入、权益兑付、分润与品牌授权尚待信用卡中心论证"),
        ("宣传口径", "联合品牌/联名卡表述须授权并经宣传与卡中心审核"),
    ]
    for i, (t, d) in enumerate(risks):
        y = Inches(1.55 + i * 0.95)
        add_round(s, Inches(0.45), y, Inches(6.3), Inches(0.88), SOFT)
        add_text(s, Inches(0.65), y + Inches(0.1), Inches(5.9), Inches(0.28), t, size=13, bold=True, color=NAVY)
        add_text(s, Inches(0.65), y + Inches(0.4), Inches(5.9), Inches(0.42), d, size=11, color=DARK)

    # 右侧共识
    add_text(s, Inches(7.1), Inches(1.2), Inches(5.7), Inches(0.35), "共识 / 补充方案分层", size=16, bold=True, color=NAVY)
    add_round(s, Inches(7.0), Inches(1.55), Inches(5.85), Inches(5.2), SOFT)
    add_text(s, Inches(7.25), Inches(1.7), Inches(5.4), Inches(0.3), "已达成共识", size=13, bold=True, color=GOLD)
    consensus = [
        "优先：积分商城兑换 + 开户赠礼",
        "兑换码/License 客户自助激活",
        "套餐与阶梯价；腾讯交完整方案",
        "8/5与分行领导沟通；上海试点",
    ]
    add_bullets(s, Inches(7.25), Inches(2.05), Inches(5.4), Inches(2.0), consensus, size=12, color=DARK, spacing=1.3)
    add_text(s, Inches(7.25), Inches(4.15), Inches(5.4), Inches(0.3), "初步方案（本次补充）", size=13, bold=True, color=NAVY)
    prelim = [
        "单列员工份额：行员采购+行外采购",
        "行内：办公提效系统/小程序",
        "行内：员工AI福利发放",
        "对外：研究联名信用卡方案",
    ]
    add_bullets(s, Inches(7.25), Inches(4.5), Inches(5.4), Inches(2.0), prelim, size=12, color=DARK, spacing=1.3)
    footer(s, 4)

    # ========== 第5页：请示事项与下一步 ==========
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "四、建议请示事项与下一步", "原则同意推进 · 补充员工份额 · 研究联名卡 · 先试点后扩量")

    decisions = [
        ("01", "原则同意继续推进", "客户增值服务 + 行内数字化提效继续论证"),
        ("02", "同意上海地区试点", "选定网点/客群，设定兑换激活获客指标"),
        ("03", "补充员工采购份额", "行员采购+行外采购；福利可先行，系统须评估"),
        ("04", "研究联名信用卡", "授权卡中心论证，不作为立即发卡承诺"),
        ("05", "跨部门评估+继续谈判", "1,500万为测算框架；首期金额另行审批"),
    ]
    # 前4项2x2，第5项通栏
    for i, (num, title, desc) in enumerate(decisions[:4]):
        x = Inches(0.4 + (i % 2) * 6.45)
        y = Inches(1.2 + (i // 2) * 1.35)
        add_round(s, x, y, Inches(6.2), Inches(1.2), SOFT)
        add_round(s, x + Inches(0.15), y + Inches(0.2), Inches(0.7), Inches(0.75), NAVY)
        add_text(
            s,
            x + Inches(0.15),
            y + Inches(0.2),
            Inches(0.7),
            Inches(0.75),
            num,
            size=16,
            bold=True,
            color=GOLD,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(s, x + Inches(1.05), y + Inches(0.22), Inches(4.9), Inches(0.4), title, size=15, bold=True, color=NAVY)
        add_text(s, x + Inches(1.05), y + Inches(0.65), Inches(4.9), Inches(0.4), desc, size=12, color=DARK)

    num, title, desc = decisions[4]
    add_round(s, Inches(0.4), Inches(3.95), Inches(12.5), Inches(0.95), SOFT)
    add_round(s, Inches(0.55), Inches(4.05), Inches(0.7), Inches(0.75), NAVY)
    add_text(
        s,
        Inches(0.55),
        Inches(4.05),
        Inches(0.7),
        Inches(0.75),
        num,
        size=16,
        bold=True,
        color=GOLD,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(s, Inches(1.45), Inches(4.1), Inches(11.2), Inches(0.35), title, size=15, bold=True, color=NAVY)
    add_text(s, Inches(1.45), Inches(4.5), Inches(11.2), Inches(0.35), desc, size=12, color=DARK)

    # 底部结论条
    add_round(s, Inches(0.4), Inches(5.1), Inches(12.5), Inches(1.7), NAVY)
    add_text(s, Inches(0.7), Inches(5.25), Inches(11.9), Inches(0.3), "汇报结论 / 下一步", size=13, bold=True, color=GOLD)
    add_text(
        s,
        Inches(0.7),
        Inches(5.6),
        Inches(11.9),
        Inches(1.0),
        "建议原则同意立项论证并授权继续谈判，同步论证员工份额与联名信用卡；先试点、后扩量。\n"
        "下一步：① 8/5前腾讯交完整方案（含行外/行员份额）　② 300万/500万两档　③ 兑换激活演示\n"
        "　　　　④ 内部系统合规边界　⑤ 联名卡预研　⑥ 跨部门预审与试点指标　⑦ 达标后再扩量采购",
        size=12,
        color=WHITE,
    )
    footer(s, 5)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"已生成：{output_path}")


if __name__ == "__main__":
    out = Path(__file__).resolve().parents[1] / "deliverables" / "泰隆银行上海分行与腾讯WorkBuddy合作推进情况汇报_20260803.pptx"
    build_ppt(out)
