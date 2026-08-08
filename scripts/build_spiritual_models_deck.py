#!/usr/bin/env python3
"""从一场灵性闲聊总结中，形式化推导可演示的概念模型，并生成 PPTX。

素材来源：2026-08-08 线下多人闲聊智能总结（约 61 分钟）。
本脚本不做宗教真实性判断，只做「民间本体论 → 可操作模型」的结构化抽取。
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT = Path("/workspace/outputs/灵性闲聊_可推导模型.pptx")

# 视觉方向：墨青学术风 —— 避开紫渐变 / 奶油衬线赤陶 / 报纸排版
INK = RGBColor(0x14, 0x28, 0x24)       # 深墨青
INK2 = RGBColor(0x1E, 0x3A, 0x34)      # 次墨
CELADON = RGBColor(0x4F, 0x8A, 0x78)   # 青瓷
AMBER = RGBColor(0xC4, 0xA3, 0x5A)     # 暖琥珀点缀
MIST = RGBColor(0xE8, 0xEF, 0xEC)      # 雾青底
PANEL = RGBColor(0xF4, 0xF7, 0xF5)     # 面板
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1C, 0x1C, 0x1C)
MUTED = RGBColor(0x5A, 0x6B, 0x66)
SOFT = RGBColor(0x2F, 0x4A, 0x44)

FONT = "Microsoft YaHei"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, *, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def add_round(slide, x, y, w, h, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    shp.adjustments[0] = 0.08
    return shp


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    size=18,
    bold=False,
    color=DARK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=15, color=DARK, bullet=CELADON):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.28
        r1 = p.add_run()
        r1.text = "●  "
        r1.font.name = FONT
        r1.font.size = Pt(size)
        r1.font.color.rgb = bullet
        r2 = p.add_run()
        r2.text = item
        r2.font.name = FONT
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
    return tb


def header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SW, Inches(0.92), INK)
    add_rect(slide, 0, Inches(0.92), SW, Inches(0.05), CELADON)
    add_text(slide, Inches(0.55), Inches(0.18), Inches(12), Inches(0.45), title, size=24, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.55), Inches(0.55), Inches(12), Inches(0.32), subtitle, size=12, color=MIST)


def footer(slide, page, total=12):
    add_text(
        slide,
        Inches(0.55),
        Inches(7.12),
        Inches(10),
        Inches(0.28),
        "灵性闲聊 · 可推导模型 · 2026-08-08 录音总结二次抽取",
        size=10,
        color=MUTED,
    )
    add_text(
        slide,
        Inches(11.6),
        Inches(7.12),
        Inches(1.2),
        Inches(0.28),
        f"{page}/{total}",
        size=10,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def card(slide, x, y, w, h, title, body_lines, *, accent=CELADON):
    add_round(slide, x, y, w, h, WHITE)
    add_rect(slide, x, y, Inches(0.08), h, accent)
    add_text(slide, x + Inches(0.25), y + Inches(0.18), w - Inches(0.4), Inches(0.35), title, size=16, bold=True, color=INK)
    add_bullets(slide, x + Inches(0.25), y + Inches(0.55), w - Inches(0.4), h - Inches(0.7), body_lines, size=13)


# ========== 1 封面 ==========
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, INK)
# 大气层：斜向色块营造空间感
add_rect(s, Inches(8.8), 0, Inches(4.6), SH, INK2)
add_rect(s, 0, Inches(5.85), SW, Inches(0.04), CELADON)
add_rect(s, Inches(0.8), Inches(1.35), Inches(1.4), Inches(0.05), AMBER)
add_text(s, Inches(0.8), Inches(1.55), Inches(8), Inches(0.4), "从闲聊中抽取可操作框架", size=16, color=AMBER)
add_text(s, Inches(0.8), Inches(2.15), Inches(10), Inches(1.2), "灵性闲聊的可推导模型", size=44, bold=True, color=WHITE)
add_text(
    s,
    Inches(0.8),
    Inches(3.45),
    Inches(9.5),
    Inches(0.9),
    "把「前世回溯 / 能量链接 / 佛道双修 / 职场接手」\n压缩成 6 套可对照、可迁移的概念模型",
    size=18,
    color=MIST,
)
add_text(s, Inches(0.8), Inches(5.0), Inches(9), Inches(0.4), "素材 · 2026-08-08 12:37–13:39 · 约 6 人线下闲聊 · 61 分钟", size=14, color=CELADON)
add_text(s, Inches(0.8), Inches(6.2), Inches(9), Inches(0.35), "方法：民间本体论 → 变量抽取 → 结构方程 → 跨域同构", size=13, color=MUTED)
add_text(s, Inches(9.3), Inches(6.2), Inches(3.5), Inches(0.35), "MODEL STUDIO", size=13, color=AMBER, align=PP_ALIGN.RIGHT)

# ========== 2 总览 ==========
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, MIST)
header(s, "六套模型总览", "同一场闲聊，抽出六种可并列对照的结构")
footer(s, 2)

models = [
    ("M1 能量存量—流量", "链接是否耗能，取决于自身存量与链接维度高低"),
    ("M2 输出配额与破例", "常规上限三次；叙事无聊可触发第四次破例"),
    ("M3 跨世缘分图", "促成婚配 → 离散 → 重逢生育 → 私塾教学的角色连续"),
    ("M4 非线性并行时空", "时间非线性；另一时空中的「我」可能同步在工作"),
    ("M5 佛道双修互补栈", "道重护身延寿，佛需能量与认知门槛，二者不互斥"),
    ("M6 接手修复循环", "多任失败 → 信心崩塌 → 全量验货 → 重建 → 转向"),
]
for i, (t, d) in enumerate(models):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(4.2) * col
    y = Inches(1.25) + Inches(2.7) * row
    add_round(s, x, y, Inches(4.0), Inches(2.4), WHITE)
    add_rect(s, x, y, Inches(4.0), Inches(0.08), CELADON if i % 2 == 0 else AMBER)
    add_text(s, x + Inches(0.25), y + Inches(0.35), Inches(3.5), Inches(0.45), t, size=18, bold=True, color=INK)
    add_text(s, x + Inches(0.25), y + Inches(1.0), Inches(3.5), Inches(1.1), d, size=14, color=SOFT)

# ========== 3 M1 能量 ==========
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, MIST)
header(s, "M1 · 能量存量—流量模型", "核心金句：「能量消耗与否要看本身存量；链接高位则每次补充」")
footer(s, 3)

add_round(s, Inches(0.5), Inches(1.25), Inches(7.4), Inches(5.5), WHITE)
add_text(s, Inches(0.8), Inches(1.45), Inches(6.8), Inches(0.4), "形式化表达", size=16, bold=True, color=INK)
add_text(
    s,
    Inches(0.8),
    Inches(1.95),
    Inches(6.8),
    Inches(2.2),
    "设个体能量存量为 S，链接对象维度差为 ΔD。\n\n"
    "ΔS = g(S, ΔD)\n\n"
    "· 若 ΔD > 0（链接更高维能量位）：ΔS ≥ 0，每次补充\n"
    "· 若 ΔD ≤ 0：是否消耗取决于 S 是否低于阈值 θ\n"
    "· 低存量 + 平级/低位链接 → 净消耗",
    size=15,
    color=DARK,
)
add_bullets(
    s,
    Inches(0.8),
    Inches(4.3),
    Inches(6.8),
    Inches(2.2),
    [
        "可迁移：社交、信息流、上下级沟通同样适用「存量 × 维度差」",
        "诊断问题从「要不要断联」变为「我的 S 是否过低、对方是否高位」",
        "高位链接 = 稀缺资源；平级闲聊在低 S 时成本最高",
    ],
    size=14,
)

# 右侧决策卡
add_round(s, Inches(8.15), Inches(1.25), Inches(4.65), Inches(5.5), INK)
add_text(s, Inches(8.45), Inches(1.55), Inches(4.1), Inches(0.4), "决策矩阵", size=16, bold=True, color=AMBER)
rows = [
    ("高 S × 高位", "净增益 · 可加深"),
    ("高 S × 平级", "近似中性 · 可控"),
    ("低 S × 高位", "优先补充 · 救命"),
    ("低 S × 平级/低位", "易透支 · 应节流"),
]
for i, (k, v) in enumerate(rows):
    y = Inches(2.2) + Inches(1.0) * i
    add_round(s, Inches(8.45), y, Inches(4.1), Inches(0.85), INK2)
    add_text(s, Inches(8.65), y + Inches(0.12), Inches(3.7), Inches(0.3), k, size=13, bold=True, color=CELADON)
    add_text(s, Inches(8.65), y + Inches(0.42), Inches(3.7), Inches(0.3), v, size=13, color=WHITE)

# ========== 4 M2 配额 ==========
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, MIST)
header(s, "M2 · 灵性输出配额与「无聊破例」", "常规上限 3 次；第 4 次因「剧本无聊想更换」而破例")
footer(s, 4)

card(
    s,
    Inches(0.5),
    Inches(1.25),
    Inches(6.0),
    Inches(5.5),
    "规则层",
    [
        "每人此类输出默认硬预算 n_max = 3",
        "约束来源：能量不够（资源约束，非道德约束）",
        "本次已是第 4 次 = 预算外特例",
        "触发条件：讲述者对既有「剧本」感到无聊",
        "破例功能：刷新叙事场景，而非单纯加量",
        "新场景仍保持气质连续（高大威猛、气质一致）",
    ],
)
card(
    s,
    Inches(6.8),
    Inches(1.25),
    Inches(6.0),
    Inches(5.5),
    "有意思的推论",
    [
        "能量预算可被「叙事质量」推翻",
        "无聊 = 一种系统级纠错信号",
        "破例不是无限透支，而是换剧本降熵",
        "同构于产品：第三版无聊 → 允许第四版重构",
        "同构于研究：重复实验无新意 → 允许改设问",
        "关键指标：不是次数本身，而是边际新颖度",
    ],
    accent=AMBER,
)

# ========== 5 M3 缘分图 ==========
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, MIST)
header(s, "M3 · 跨世缘分图（角色连续）", "两段古代人生 + 面相读前世 → 今生教育属性匹配")
footer(s, 5)

# 时间线节点
nodes = [
    ("世 1", "促成男生结婚\n之后分散离开"),
    ("世 2", "家道中落姑娘\n重逢 · 多子女\n男生开私塾"),
    ("面相", "书卷气重\n适合教育业"),
    ("今生", "工作属性\n类似老师\n匹配度高"),
]
for i, (t, d) in enumerate(nodes):
    x = Inches(0.55) + Inches(3.15) * i
    add_round(s, x, Inches(1.4), Inches(2.9), Inches(3.2), WHITE)
    add_rect(s, x, Inches(1.4), Inches(2.9), Inches(0.55), CELADON if i < 3 else AMBER)
    add_text(s, x, Inches(1.48), Inches(2.9), Inches(0.4), t, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(2.2), Inches(2.5), Inches(2.1), d, size=14, color=DARK, align=PP_ALIGN.CENTER)
    if i < 3:
        add_text(s, x + Inches(2.85), Inches(2.6), Inches(0.4), Inches(0.4), "→", size=22, bold=True, color=CELADON)

add_round(s, Inches(0.55), Inches(4.9), Inches(12.2), Inches(1.8), WHITE)
add_text(s, Inches(0.85), Inches(5.1), Inches(11.6), Inches(0.35), "结构洞见：气质连续 + 功能连续", size=15, bold=True, color=INK)
add_text(
    s,
    Inches(0.85),
    Inches(5.55),
    Inches(11.6),
    Inches(0.9),
    "故事核不是「前世浪漫」，而是「教育者角色跨世复现」：促成婚配者 → 私塾配偶场域 → 面相书卷气 → 今生类教师工作。\n"
    "可抽象为：PastRole ≈ FaceFeature ≈ PresentOccupation。匹配成功被当作验证信号。",
    size=14,
    color=SOFT,
)

# ========== 6 M4 时空 ==========
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, MIST)
header(s, "M4 · 非线性并行时空", "「时间不是线性的，时空是伪的」")
footer(s, 6)

add_round(s, Inches(0.5), Inches(1.25), Inches(4.0), Inches(5.5), INK)
add_text(s, Inches(0.8), Inches(1.55), Inches(3.4), Inches(0.4), "命题", size=16, bold=True, color=AMBER)
add_bullets(
    s,
    Inches(0.8),
    Inches(2.2),
    Inches(3.4),
    Inches(4.2),
    [
        "时间非线性",
        "时空概念本身虚假",
        "另一时空中，在场三人可能同步在工作",
        "当下自我只是多投影之一",
        "「还没吃饭」与「彼岸在上班」可并存",
    ],
    size=14,
    color=WHITE,
    bullet=AMBER,
)

add_round(s, Inches(4.8), Inches(1.25), Inches(7.9), Inches(5.5), WHITE)
add_text(s, Inches(5.15), Inches(1.55), Inches(7.2), Inches(0.4), "可推导的有趣结构", size=16, bold=True, color=INK)
add_bullets(
    s,
    Inches(5.15),
    Inches(2.2),
    Inches(7.2),
    Inches(4.2),
    [
        "并行自我 = 多线程进程：本地闲聊线程与远端工作线程同时运行",
        "「伪时空」≈ 把线性时钟降级为本地坐标系，而非本体论实在",
        "与 M1 联立：高维链接 = 跨线程 RPC，可能回灌能量",
        "与 M2 联立：换剧本 = 切换到另一条叙事时间线",
        "与职场故事同构：接手失败项目 = 切入被放弃的时间线分支并修复",
        "实用隐喻：决策时假设「另一版本的我已在推进」，降低拖延焦虑",
    ],
    size=14,
)

# ========== 7 M5 佛道 ==========
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, MIST)
header(s, "M5 · 佛道双修互补栈", "各有精华，不必互相排斥")
footer(s, 7)

card(
    s,
    Inches(0.5),
    Inches(1.25),
    Inches(6.0),
    Inches(3.6),
    "道家层 · 护身与延寿",
    [
        "护身功法被评价为「非常好」",
        "功能定位：身体保全、延年益寿",
        "偏工程：可练习、可体感、可日常化",
        "适合作为底座层 / 操作系统层",
    ],
)
card(
    s,
    Inches(6.8),
    Inches(1.25),
    Inches(6.0),
    Inches(3.6),
    "佛家层 · 认知门槛",
    [
        "金刚经等经典难真正看懂",
        "障碍 = 能量不足 + 认知基础不足",
        "偏理论：高门槛应用层",
        "需先有底座能量，再进入深度文本",
    ],
    accent=AMBER,
)

add_round(s, Inches(0.5), Inches(5.1), Inches(12.3), Inches(1.6), WHITE)
add_text(s, Inches(0.8), Inches(5.3), Inches(11.8), Inches(0.35), "栈式理解：先护身（道）再读经（佛）", size=15, bold=True, color=INK)
add_text(
    s,
    Inches(0.8),
    Inches(5.75),
    Inches(11.8),
    Inches(0.7),
    "这不是宗教裁判，而是学习路径模型：低层稳定生理/能量，高层才承载抽象经典。\n"
    "同构于技术栈：先把基础设施跑稳，再上复杂业务语义。",
    size=14,
    color=SOFT,
)

# ========== 8 M6 职场 ==========
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, MIST)
header(s, "M6 · 失败项目接手修复循环", "多任失败 → 信心崩塌 → 全量验货 → 按自己的思路推进 → 转向")
footer(s, 8)

steps = [
    ("1 崩塌", "多任项目组\n做不成\n团队失信"),
    ("2 接手", "项目转交\n责任下沉\n到讲述者"),
    ("3 验货", "好产品\n坏产品\n全部查验"),
    ("4 重建", "按自己\n思路推进\n完成交付"),
    ("5 转向", "进入新群体\n选适合的路\n换发展方向"),
]
for i, (t, d) in enumerate(steps):
    x = Inches(0.45) + Inches(2.55) * i
    add_round(s, x, Inches(1.35), Inches(2.35), Inches(3.4), WHITE if i != 2 else INK)
    title_c = INK if i != 2 else AMBER
    body_c = DARK if i != 2 else WHITE
    add_text(s, x + Inches(0.15), Inches(1.55), Inches(2.05), Inches(0.45), t, size=16, bold=True, color=title_c, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), Inches(2.2), Inches(2.05), Inches(2.2), d, size=14, color=body_c, align=PP_ALIGN.CENTER)

add_round(s, Inches(0.45), Inches(5.05), Inches(12.4), Inches(1.65), WHITE)
add_text(s, Inches(0.75), Inches(5.25), Inches(11.9), Inches(0.35), "与灵性模型的同构", size=15, bold=True, color=INK)
add_text(
    s,
    Inches(0.75),
    Inches(5.7),
    Inches(11.9),
    Inches(0.75),
    "「全量验货」≈ 对好坏剧本/好坏业力不做过滤；「换思路」≈ M2 换剧本；「转向新群体」≈ 切到更高维场域。\n"
    "关键动作不是鸡汤鼓励，而是先拿到完整信息集，再重建局部秩序。",
    size=14,
    color=SOFT,
)

# ========== 9 身份场域 ==========
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, MIST)
header(s, "附加模型 · 身份场域与「命名先于空间」", "从自我介绍与 21 层装修闲聊中抽出")
footer(s, 9)

card(
    s,
    Inches(0.5),
    Inches(1.25),
    Inches(6.0),
    Inches(5.5),
    "身份三角",
    [
        "地理锚点：复旦大学周边专职工作",
        "组织锚点：杨浦区科技企业联合中心秘书长",
        "业务三角：新级增长 / 不动产 / 人工智能",
        "灵性意识被定义为「类似小孩的敏感度」",
        "在场评价：对话对象灵性意识状态不错",
        "推论：专业身份与灵性场域并行不冲突",
    ],
)
card(
    s,
    Inches(6.8),
    Inches(1.25),
    Inches(6.0),
    Inches(5.5),
    "命名先于楼层",
    [
        "先定名称「21」，后选 21 层场地",
        "意向 → 命名 → 物理落位",
        "材料：21 层及以上统一 10 斤规格",
        "可做顶光；上下规格一致",
        "隐喻：品牌/数字先于选址",
        "同构于创业：先定符号，再找承载空间",
    ],
    accent=AMBER,
)

# ========== 10 总联立 ==========
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, MIST)
header(s, "六模型联立：一张「闲聊操作系统」", "把松散话题压成可运行的心智 OS")
footer(s, 10)

add_round(s, Inches(0.5), Inches(1.25), Inches(12.3), Inches(5.5), WHITE)
layers = [
    ("感知层", "小孩式敏感度 · 面相/气质读取 · 当场匹配验证"),
    ("叙事层", "跨世缘分图 · 剧本质量监控 · 无聊则换剧本"),
    ("能量层", "存量 S · 维度差 ΔD · 高位补充 / 低位节流"),
    ("修持层", "道：护身延寿底座 · 佛：高门槛语义层"),
    ("行动层", "全量验货接手 · 重建局部秩序 · 完成后转向"),
    ("场域层", "命名先于空间 · 身份三角并行 · 物理规格统一"),
]
for i, (t, d) in enumerate(layers):
    y = Inches(1.5) + Inches(0.8) * i
    add_rect(s, Inches(0.85), y, Inches(2.4), Inches(0.65), CELADON if i % 2 == 0 else INK)
    add_text(s, Inches(0.85), y, Inches(2.4), Inches(0.65), t, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(3.5), y, Inches(8.8), Inches(0.65), d, size=15, color=DARK, anchor=MSO_ANCHOR.MIDDLE)

# ========== 11 可迁移玩法 ==========
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, MIST)
header(s, "有意思的迁移：把它当成设计工具", "不争论真假，先榨取结构价值")
footer(s, 11)

ideas = [
    ("产品设计", "给功能设「三次常规输出」；第四次必须换剧本才允许"),
    ("团队管理", "接手烂摊子先全量验货，禁止只看成功样例"),
    ("个人精力", "用 S×ΔD 矩阵管理社交与信息输入"),
    ("学习路径", "先建护身/作息底座，再碰高门槛经典或理论"),
    ("品牌选址", "先锁定符号与名称，再反推物理空间"),
    ("叙事疗愈", "把重复故事标为「无聊」，主动触发场景刷新"),
]
for i, (t, d) in enumerate(ideas):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(4.2) * col
    y = Inches(1.3) + Inches(2.7) * row
    add_round(s, x, y, Inches(4.0), Inches(2.4), WHITE)
    add_text(s, x + Inches(0.25), y + Inches(0.35), Inches(3.5), Inches(0.4), t, size=18, bold=True, color=INK)
    add_text(s, x + Inches(0.25), y + Inches(1.0), Inches(3.5), Inches(1.1), d, size=14, color=SOFT)

# ========== 12 结尾 ==========
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, INK)
add_rect(s, 0, Inches(5.9), SW, Inches(0.04), CELADON)
add_text(s, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.5), "一句话收束", size=16, color=AMBER)
add_text(
    s,
    Inches(0.8),
    Inches(2.4),
    Inches(11.5),
    Inches(1.4),
    "这场闲聊表面松散，内里其实在谈：\n如何管理能量、刷新叙事、跨世/跨项目修复，以及如何把传统双修当成学习栈。",
    size=22,
    bold=True,
    color=WHITE,
)
add_text(
    s,
    Inches(0.8),
    Inches(4.2),
    Inches(11.5),
    Inches(1.0),
    "最有意思的不是「信或不信」，而是这些民间说法已经自带变量、阈值、破例条件和跨域同构。",
    size=16,
    color=MIST,
)
add_text(s, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.4), "交付物 · outputs/灵性闲聊_可推导模型.pptx · 脚本 scripts/build_spiritual_models_deck.py", size=12, color=MUTED)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"已生成: {OUT}")
