#!/usr/bin/env python3
"""生成《发展与去留：2026 中美判断备忘录》PPT。"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

INK = RGBColor(0x0F, 0x17, 0x2A)
INK_SOFT = RGBColor(0x1E, 0x29, 0x3B)
TEAL = RGBColor(0x0F, 0x76, 0x6E)
TEAL_LIGHT = RGBColor(0x14, 0xB8, 0xA6)
AMBER = RGBColor(0xD9, 0x77, 0x06)
ROSE = RGBColor(0xBE, 0x12, 0x31)
SLATE = RGBColor(0x47, 0x55, 0x69)
SOFT = RGBColor(0xF1, 0xF5, 0xF9)
MINT = RGBColor(0xEC, 0xFD, 0xF5)
CREAM = RGBColor(0xFF, 0xF7, 0xED)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x29, 0x3B)
LINE = RGBColor(0xCB, 0xD5, 0xE1)

FONT = "Microsoft YaHei"
FOOTER = "发展与去留 · 2026 中美判断备忘录 · 默认深耕，出海做成期权"
TOTAL = 16
SW = SH = None


def add_rect(slide, x, y, w, h, fill, *, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_round(slide, x, y, w, h, fill, *, adj=0.08):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    shp.adjustments[0] = adj
    return shp


def set_run(run, text, *, size=16, bold=False, color=DARK, font=FONT):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    size=16,
    bold=False,
    color=DARK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        set_run(p.add_run(), line, size=size, bold=bold, color=color)
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=14, color=DARK, mark="●", mark_color=TEAL, spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Emu(0))
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = spacing
        set_run(p.add_run(), f"{mark}  ", size=size, color=mark_color)
        set_run(p.add_run(), item, size=size, color=color)
    return tb


def header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SW, Inches(0.9), INK)
    add_rect(slide, 0, Inches(0.9), SW, Inches(0.06), TEAL_LIGHT)
    add_text(slide, Inches(0.5), Inches(0.18), Inches(12.2), Inches(0.4), title, size=22, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.52), Inches(12.2), Inches(0.3), subtitle, size=12, color=TEAL_LIGHT)


def footer(slide, page):
    add_text(slide, Inches(0.5), Inches(7.12), Inches(10.8), Inches(0.26), FOOTER, size=10, color=SLATE)
    add_text(
        slide,
        Inches(11.4),
        Inches(7.12),
        Inches(1.4),
        Inches(0.26),
        f"{page} / {TOTAL}",
        size=10,
        color=SLATE,
        align=PP_ALIGN.RIGHT,
    )


def card(slide, x, y, w, h, title, items, *, accent=TEAL, bg=SOFT, title_size=14, item_size=12):
    add_round(slide, x, y, w, h, bg)
    add_rect(slide, x, y, Inches(0.09), h, accent)
    add_text(
        slide,
        x + Inches(0.26),
        y + Inches(0.14),
        w - Inches(0.4),
        Inches(0.34),
        title,
        size=title_size,
        bold=True,
        color=accent,
    )
    add_bullets(
        slide,
        x + Inches(0.26),
        y + Inches(0.5),
        w - Inches(0.42),
        h - Inches(0.62),
        items,
        size=item_size,
        mark_color=accent,
    )


def kpi(slide, x, y, w, h, num, label, note, *, accent=TEAL):
    add_round(slide, x, y, w, h, SOFT)
    add_text(slide, x + Inches(0.16), y + Inches(0.12), w - Inches(0.3), Inches(0.42), num, size=22, bold=True, color=accent)
    add_text(slide, x + Inches(0.16), y + Inches(0.54), w - Inches(0.3), Inches(0.28), label, size=12, bold=True, color=INK)
    add_text(slide, x + Inches(0.16), y + Inches(0.82), w - Inches(0.3), Inches(0.36), note, size=11, color=SLATE)


def build_ppt(output_path: Path) -> None:
    global SW, SH
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    # 1 封面
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, INK)
    add_rect(s, 0, 0, Inches(0.18), SH, TEAL_LIGHT)
    add_text(s, Inches(0.8), Inches(1.15), Inches(11.5), Inches(0.32), "个人判断备忘录 · 2026 年 8 月", size=14, color=TEAL_LIGHT)
    add_text(s, Inches(0.8), Inches(1.6), Inches(11.8), Inches(1.35), "发展与去留", size=48, bold=True, color=WHITE)
    add_text(
        s,
        Inches(0.8),
        Inches(3.0),
        Inches(11.5),
        Inches(0.7),
        "现在的中国，适不适合继续发展？\n人工智能的机会，是就地抓住，还是果断去海外？",
        size=18,
        color=SOFT,
    )
    add_rect(s, Inches(0.8), Inches(3.9), Inches(2.2), Inches(0.06), TEAL_LIGHT)
    add_text(
        s,
        Inches(0.8),
        Inches(4.2),
        Inches(11.5),
        Inches(0.7),
        "结论先行：对网络型、场景型、转化型从业者，默认留在中国深耕；\n把「去美国」做成有限期期权，而不是斩根式流亡。",
        size=16,
        color=TEAL_LIGHT,
    )
    pills = [("01", "先破二元"), ("02", "再拆那段话"), ("03", "后定路径")]
    for i, (num, name) in enumerate(pills):
        x = Inches(0.8 + i * 3.6)
        add_round(s, x, Inches(5.45), Inches(3.3), Inches(0.95), INK_SOFT)
        add_text(s, x + Inches(0.22), Inches(5.58), Inches(0.8), Inches(0.7), num, size=18, bold=True, color=TEAL_LIGHT, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(1.05), Inches(5.58), Inches(2.05), Inches(0.7), name, size=18, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    # 2 三问
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "先把问题拆开", "对方给的是「如何走」的诀别辞，你要回答的是「应不应该走」")
    qs = [
        ("问一", "留在中国，还是去美国？", "这是地理选择题，但真正要选的是能力结构放在哪边更值钱。"),
        ("问二", "现在的中国，还适不适合继续发展？", "「发展」必须定义：是杠杆地产，还是新动能上的场景与转化。"),
        ("问三", "AI 机会能不能抓住？", "抓住的是应用、产业空间与商业化，还是美西少数实验室的前沿训练。"),
    ]
    for i, (tag, title, note) in enumerate(qs):
        y = Inches(1.2 + i * 1.55)
        add_round(s, Inches(0.5), y, Inches(12.3), Inches(1.4), SOFT)
        add_rect(s, Inches(0.5), y, Inches(0.12), Inches(1.4), TEAL if i != 1 else AMBER)
        add_text(s, Inches(0.9), y + Inches(0.18), Inches(1.2), Inches(0.32), tag, size=13, bold=True, color=TEAL)
        add_text(s, Inches(2.1), y + Inches(0.16), Inches(10.3), Inches(0.4), title, size=20, bold=True, color=INK)
        add_text(s, Inches(0.9), y + Inches(0.72), Inches(11.5), Inches(0.5), note, size=14, color=SLATE)
    footer(s, 2)

    # 3 结论先行
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "结论先行：不是「中国完了」，也不是「必须留下」", "对你这种能力结构，2026 年的最优解是深耕 + 期权，不是斩根")
    card(
        s,
        Inches(0.45),
        Inches(1.2),
        Inches(4.0),
        Inches(5.5),
        "适合留下的理由",
        [
            "中国仍是全球最大的 AI 应用与制造业融合市场",
            "新动能对上半年增长贡献已超四成",
            "你的资产是本地网络、场景与转化，搬不走",
            "上海 / 杨浦 / 高校 / 产业空间仍是可操作战场",
        ],
        accent=TEAL,
        bg=MINT,
    )
    card(
        s,
        Inches(4.65),
        Inches(1.2),
        Inches(4.0),
        Inches(5.5),
        "不适合幻想的部分",
        [
            "地产开发投资仍在深度调整，旧剧本回不来",
            "民间投资与预期偏弱，不是 2015 年的顺风",
            "Frontier 大模型资本与人才密度仍在美西",
            "「再熬一波房价」不是发展策略",
        ],
        accent=AMBER,
        bg=CREAM,
    )
    card(
        s,
        Inches(8.85),
        Inches(1.2),
        Inches(4.0),
        Inches(5.5),
        "去美国的真实约束",
        [
            "2026 年高技能移民通道显著收紧",
            "H-1B 成本与不确定性上升，身份即风险",
            "没有具体岗位 / 签证 / 薪酬溢价，就是裸奔",
            "中年网络型从业者，美国是重置而不是升级",
        ],
        accent=ROSE,
        bg=SOFT,
    )
    footer(s, 3)

    # 4 那段话
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "如何看待那段话", "它很美，也很危险：美在给行动力，险在把乡愁问题当成了战略答案")
    add_round(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.15), INK)
    add_text(
        s,
        Inches(0.8),
        Inches(1.4),
        Inches(11.7),
        Inches(1.75),
        "「不要回来，不要想念我们，不要回头……不要落叶归根，而要落地生根。\n心安之处即是家。别回头，向前走。不归根，只生根。心安处，即吾乡。」",
        size=16,
        color=WHITE,
    )
    card(
        s,
        Inches(0.5),
        Inches(3.55),
        Inches(6.05),
        Inches(3.15),
        "它真正解决的问题",
        [
            "移民史上常见的「斩断嘱咐」：怕乡愁把人撕成两半",
            "给远行者一个操作系统：把思念从决策权里拿掉",
            "半句内核成立——心安处即吾乡，家不必等于出生地",
        ],
        accent=TEAL,
    )
    card(
        s,
        Inches(6.75),
        Inches(3.55),
        Inches(6.05),
        Inches(3.15),
        "它没有回答的问题",
        [
            "走，是已经决定的事，还是正在犹豫的事？",
            "封存「这里的一切」之后，你还剩什么可交易的能力？",
            "心安如果要靠切断记忆来换，那还叫不叫心安？",
        ],
        accent=AMBER,
        bg=CREAM,
    )
    footer(s, 4)

    # 5 三处误用
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "那段话的三处误用", "把「如何走」的心理技术，误当成「应不应该走」的结论")
    misuses = [
        ("误用一", "把诀别辞当成战略", "「不要回来」是走了之后的自我保护，不是走之前的可行性研究。先问值不值得走，再谈如何不回头。"),
        ("误用二", "把家简化成地理", "家当然可以重建。但职业资本、信任网络、语言与意义感都有路径依赖。换坐标不等于换命运。"),
        ("误用三", "把封存当成勇敢", "对实验室科学家，封存故土或许轻；对靠本地关系做转化的人，「封存这里的一切」等于烧掉资产负债表。"),
    ]
    for i, (tag, title, body) in enumerate(misuses):
        x = Inches(0.45 + i * 4.25)
        add_round(s, x, Inches(1.25), Inches(4.05), Inches(5.4), SOFT)
        add_rect(s, x, Inches(1.25), Inches(4.05), Inches(0.1), ROSE if i == 2 else AMBER)
        add_text(s, x + Inches(0.25), Inches(1.55), Inches(3.55), Inches(0.35), tag, size=13, bold=True, color=ROSE)
        add_text(s, x + Inches(0.25), Inches(2.0), Inches(3.55), Inches(0.9), title, size=22, bold=True, color=INK)
        add_text(s, x + Inches(0.25), Inches(3.05), Inches(3.55), Inches(3.2), body, size=14, color=SLATE)
    footer(s, 5)

    # 6 双速经济
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "2026 上半年中国：不是崩盘，是双速", "总量仍在合理区间，结构已经换挡；旧引擎冷、新动能热")
    kpis = [
        ("+4.7%", "GDP 同比增长", "上半年 69.57 万亿；Q2 回落到 4.3%"),
        ("5.2%", "城镇调查失业率", "均值与上年持平，就业总体稳定"),
        ("-18.0%", "房地产开发投资", "销售面积 -11.6%，旧引擎仍在去杠杆"),
        (">40%", "新动能贡献率", "高技术制造 +13.3%，电子利润 +96.9%"),
    ]
    for i, (num, label, note) in enumerate(kpis):
        x = Inches(0.45 + i * 3.2)
        accent = ROSE if i == 2 else TEAL
        kpi(s, x, Inches(1.2), Inches(3.05), Inches(1.3), num, label, note, accent=accent)
    add_text(s, Inches(0.5), Inches(2.7), Inches(12.3), Inches(0.35), "读法：宏观「还活着」不等于人人都能用旧方法赚钱。真正的分叉在你站在哪一条速度上。", size=14, color=SLATE)
    card(
        s,
        Inches(0.45),
        Inches(3.15),
        Inches(6.1),
        Inches(3.5),
        "冷的一侧",
        [
            "固定资产投资 -5.7%，民间投资 -8.5%",
            "房地产销售与投资继续下滑，风险仍在出清",
            "二季度增速回落，外需与预期仍不稳定",
            "若你的主业仍绑定增量开发，会越来越难",
        ],
        accent=ROSE,
    )
    card(
        s,
        Inches(6.75),
        Inches(3.15),
        Inches(6.1),
        Inches(3.5),
        "热的一侧",
        [
            "高技术制造、数字产品、装备制造明显快于整体",
            "集成电路、智能车载设备等 AI 相关行业高增长",
            "信息传输 / 软件服务业增加值 +10.7%",
            "若你做场景、空间、转化，风口在这边",
        ],
        accent=TEAL,
        bg=MINT,
    )
    footer(s, 6)

    # 7 适不适合发展
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "适不适合继续发展？取决于你做哪一种发展", "中国仍适合发展，但不适合用 2010–2020 的地产杠杆逻辑发展")
    rows = [
        ("旧发展", "拿地、加杠杆、等升值", "不适合", ROSE),
        ("守成发展", "躺平等周期、等政策、等风口回头", "不适合", AMBER),
        ("转化发展", "把 AI 接到产业、空间、政策与资本", "适合", TEAL),
        ("网络发展", "在上海做高校 / 园区 / 企业的接口人", "适合", TEAL),
    ]
    add_round(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.55), INK)
    for i, h in enumerate(["类型", "典型动作", "2026 年判断"]):
        add_text(s, Inches(0.7 + i * 4.0), Inches(1.28), Inches(3.8), Inches(0.4), h, size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    for i, (kind, action, judge, color) in enumerate(rows):
        y = Inches(1.85 + i * 1.05)
        add_round(s, Inches(0.5), y, Inches(12.3), Inches(0.95), SOFT)
        add_rect(s, Inches(0.5), y, Inches(0.1), Inches(0.95), color)
        add_text(s, Inches(0.8), y, Inches(3.6), Inches(0.95), kind, size=18, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(4.5), y, Inches(4.2), Inches(0.95), action, size=16, color=SLATE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(8.9), y, Inches(3.5), Inches(0.95), judge, size=18, bold=True, color=color, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 7)

    # 8 AI 中国
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "中国的人工智能机会：抓得住应用，抓不住幻觉", "机会真实，但不是「人人去训练大模型」；对转化者，窗口正在打开")
    card(
        s,
        Inches(0.45),
        Inches(1.2),
        Inches(6.1),
        Inches(5.5),
        "抓得住",
        [
            "制造业 + AI：设备更新、机器人、工业软件",
            "产业空间：算力、中试、园区、博物馆、会展生态",
            "商业化落地：从模型到场景的最后一公里",
            "地方政府与高校场景：可招标、可试点、可复用",
            "电子与算力需求已在利润端兑现（电子利润近翻倍）",
            "这正是「接口人」而不是「论文人」的战场",
        ],
        accent=TEAL,
        bg=MINT,
        item_size=13,
    )
    card(
        s,
        Inches(6.75),
        Inches(1.2),
        Inches(6.1),
        Inches(5.5),
        "抓不住 / 不必硬抓",
        [
            "Frontier 训练：资本、芯片、顶尖研究者仍高度集中美西",
            "没有实验室履历，去美国并不自动靠近 OpenAI 级岗位",
            "纯概念炒作、空转 PPT、无场景的「AI 平台」",
            "把 AI 当逃逸通道：地理改变不了能力结构",
            "等待「全国再来一波地产式风口」——不会再有",
            "用情绪替代订单：机会只认交付，不认焦虑",
        ],
        accent=AMBER,
        bg=CREAM,
        item_size=13,
    )
    footer(s, 8)

    # 9 美国
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "美国 2026：AI 高地还在，城门在收", "前沿仍在，但不等于对所有外国人开门；高地与关隘要分开看")
    card(
        s,
        Inches(0.45),
        Inches(1.2),
        Inches(6.1),
        Inches(3.4),
        "高地是真的",
        [
            "基础模型、算力生态、风险资本密度全球第一",
            "OpenAI / Anthropic / Nvidia 仍在抢全球顶尖人才",
            "英语世界的标准、论文、开源与薪酬锚仍在那里",
            "若你已被这类机构锁定，去，是升级不是流亡",
        ],
        accent=TEAL,
        bg=MINT,
    )
    card(
        s,
        Inches(6.75),
        Inches(1.2),
        Inches(6.1),
        Inches(3.4),
        "关隘也是真的",
        [
            "学生签、OPT、H-1B、绿卡全线收紧",
            "新 H-1B 高额费用打击中小企业与创业公司",
            "身份不确定本身会吞噬精力与议价权",
            "人才已在分流：加拿大、英国、阿联酋在抢人",
        ],
        accent=ROSE,
    )
    add_round(s, Inches(0.45), Inches(4.8), Inches(12.4), Inches(1.85), INK)
    add_text(
        s,
        Inches(0.75),
        Inches(5.05),
        Inches(11.8),
        Inches(1.4),
        "因此：「果断去海外」只对已拿到具体岗位、签证路径和薪酬溢价的人成立。\n对尚未拿到这三样的人，果断离开 = 用确定的本地网络，去换不确定的身份。这不是勇敢，是裸奔。",
        size=16,
        color=WHITE,
    )
    footer(s, 9)

    # 10 能力结构
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "关键不是中国或美国，是你的能力放在哪边更值钱", "网络、场景、转化、政策接口——这些能力有很强的地理黏性")
    add_text(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.35), "一张粗表，用来避免用别人的履历做自己的决定。", size=14, color=SLATE)
    headers_row = ["能力类型", "在中国的溢价", "在美国的溢价", "迁移成本"]
    data = [
        ["Frontier 研究 / 训练", "中（少数实验室）", "极高", "中（若已有论文与推荐）"],
        ["大厂工程 / 算法岗", "高，但内卷", "高，但签证卡脖子", "高"],
        ["产业场景转化 / 落地", "很高（政府+制造+园区）", "中低（要重做信用）", "极高"],
        ["高校 / 园区 / 商会网络", "核心资产", "接近归零", "极高"],
        ["不动产 + 产业空间", "存量时代仍有专业位", "需牌照与本地信用", "极高"],
        ["英语内容 / 跨境合作", "稀缺加分", "入场券，不是护城河", "中"],
    ]
    add_round(s, Inches(0.4), Inches(1.55), Inches(12.5), Inches(0.48), INK)
    widths = [3.1, 3.1, 3.2, 3.1]
    x0 = 0.5
    for i, h in enumerate(headers_row):
        add_text(s, Inches(x0 + sum(widths[:i])), Inches(1.58), Inches(widths[i]), Inches(0.42), h, size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    for r, row in enumerate(data):
        y = Inches(2.1 + r * 0.72)
        bg = MINT if r in (2, 3, 4) else SOFT
        add_round(s, Inches(0.4), y, Inches(12.5), Inches(0.68), bg, adj=0.04)
        for i, val in enumerate(row):
            color = TEAL if (r in (2, 3, 4) and i == 1) else INK
            add_text(
                s,
                Inches(x0 + sum(widths[:i])),
                y,
                Inches(widths[i]),
                Inches(0.68),
                val,
                size=13,
                bold=(i == 0),
                color=color,
                anchor=MSO_ANCHOR.MIDDLE,
            )
    footer(s, 10)

    # 11 决策矩阵
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "一张去留矩阵", "横轴：在本地是否还有可放大的主航道；纵轴：海外是否已有不可替代的具体机会")
    add_round(s, Inches(2.3), Inches(1.25), Inches(5.0), Inches(2.7), MINT)
    add_round(s, Inches(7.5), Inches(1.25), Inches(5.0), Inches(2.7), CREAM)
    add_round(s, Inches(2.3), Inches(4.1), Inches(5.0), Inches(2.55), SOFT)
    add_round(s, Inches(7.5), Inches(4.1), Inches(5.0), Inches(2.55), RGBColor(0xFF, 0xE4, 0xE6))
    add_text(s, Inches(2.5), Inches(1.4), Inches(4.6), Inches(0.4), "深耕中国", size=20, bold=True, color=TEAL)
    add_text(s, Inches(2.5), Inches(1.9), Inches(4.6), Inches(1.7), "本地有主航道\n海外无具体机会\n→ 2026 年大多数转化者在这里", size=14, color=SLATE)
    add_text(s, Inches(7.7), Inches(1.4), Inches(4.6), Inches(0.4), "有限期出海", size=20, bold=True, color=AMBER)
    add_text(s, Inches(7.7), Inches(1.9), Inches(4.6), Inches(1.7), "本地仍有航道\n海外也有具体机会\n→ 去 2–3 年，不斩根，预定返回条款", size=14, color=SLATE)
    add_text(s, Inches(2.5), Inches(4.25), Inches(4.6), Inches(0.4), "先造航道", size=20, bold=True, color=INK)
    add_text(s, Inches(2.5), Inches(4.75), Inches(4.6), Inches(1.6), "本地航道模糊\n海外也无机会\n→ 先不要走，先把一条业务做实", size=14, color=SLATE)
    add_text(s, Inches(7.7), Inches(4.25), Inches(4.6), Inches(0.4), "才考虑迁徙", size=20, bold=True, color=ROSE)
    add_text(s, Inches(7.7), Inches(4.75), Inches(4.6), Inches(1.6), "本地已无航道\n海外有不可替代机会\n→ 这时「果断」才有意义", size=14, color=SLATE)
    add_text(s, Inches(0.35), Inches(2.2), Inches(1.9), Inches(1.6), "海外\n已有\n机会", size=13, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.35), Inches(4.7), Inches(1.9), Inches(1.6), "海外\n尚无\n机会", size=13, bold=True, color=INK, align=PP_ALIGN.CENTER)
    footer(s, 11)

    # 12 三条路径
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "三条可执行路径", "推荐默认 A；B 是升级项；C 是纪律，不是浪漫")
    paths = [
        ("A  深耕上海", TEAL, MINT, [
            "把 AI × 产业空间 × 高校 / 园区做成一条主航道",
            "停止用情绪判断宏观，改用订单、试点、复购判断",
            "地产能力转向存量、资产管理、场景载体",
            "90 天内形成可对外讲的一条产品线",
        ]),
        ("B  双根结构", AMBER, CREAM, [
            "不去「移民」，去做 12–36 个月项目驻留或合作",
            "保留中国主体与网络，海外只做增量能力",
            "英语材料、跨境项目、对口机构，先于签证",
            "走之前写好返回条款：何时回、带回什么",
        ]),
        ("C  有条件迁徙", ROSE, SOFT, [
            "同时满足：具体岗位、签证路径、薪酬 / 意义溢价",
            "家庭与财务能承受 3 年身份不确定",
            "本地主业已无法放大，而非只是心烦",
            "即便走，也不接受「不要回来」作为人生协议",
        ]),
    ]
    for i, (title, accent, bg, items) in enumerate(paths):
        x = Inches(0.4 + i * 4.3)
        card(s, x, Inches(1.2), Inches(4.1), Inches(5.5), title, items, accent=accent, bg=bg, title_size=18, item_size=13)
    footer(s, 12)

    # 13 何时走
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "什么情况下才该走", "用清单代替口号。下列未同时点亮，就还不到「果断」的时候")
    lights = [
        ("硬条件", TEAL, ["已有雇主或机构的书面机会", "签证路径清楚，不靠赌抽签", "3 年财务缓冲，不押上全部流动性"]),
        ("事业条件", AMBER, ["去了能做这里做不了的事", "不是逃避失败，而是升级能力", "走之前本地资产有交接，而非蒸发"]),
        ("心理条件", ROSE, ["能忍受身份不确定与社会降级", "不需要靠「封存故乡」才能行动", "家人知情同意，不是一个人的诗"]),
    ]
    for i, (title, accent, items) in enumerate(lights):
        x = Inches(0.45 + i * 4.25)
        card(s, x, Inches(1.2), Inches(4.05), Inches(3.7), title, items, accent=accent, item_size=13)
    add_round(s, Inches(0.45), Inches(5.1), Inches(12.4), Inches(1.55), INK)
    add_text(
        s,
        Inches(0.75),
        Inches(5.3),
        Inches(11.8),
        Inches(1.2),
        "现在更可能的状态：本地仍有航道，海外尚无具体机会。→ 落在矩阵左上的「深耕中国」。\n那段话可以收藏，但不要让它代替这张清单。",
        size=16,
        color=WHITE,
    )
    footer(s, 13)

    # 14 90天
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "90 天行动：把判断做成行为", "前 30 天收束主航道，中 30 天验证订单，后 30 天只开一条海外期权")
    phases = [
        ("D1–30  收束", "只留一条主航道", [
            "写清：你卖的是什么，卖给谁，凭什么是你",
            "停掉所有「既要出国又要本地爆」的并行幻想",
            "把地产旧能力映射到产业空间 / 存量资产管理",
        ]),
        ("D31–60  验证", "用交付代替辩论", [
            "至少推进 1 个可报价、可试点的 AI 场景项目",
            "用周复盘看：线索、会议、订单，哪一项在动",
            "宏观新闻降权，客户反馈升权",
        ]),
        ("D61–90  期权", "出海只做增量", [
            "准备一套英文一页纸 + 项目案例",
            "锁定 1–2 个跨境合作接口（机构 / 校友 / 展会）",
            "明确：没有书面机会，就不谈移民时间表",
        ]),
    ]
    for i, (title, sub, items) in enumerate(phases):
        y = Inches(1.2 + i * 1.8)
        add_round(s, Inches(0.5), y, Inches(12.3), Inches(1.65), SOFT)
        add_rect(s, Inches(0.5), y, Inches(0.12), Inches(1.65), TEAL_LIGHT)
        add_text(s, Inches(0.9), y + Inches(0.15), Inches(3.3), Inches(0.4), title, size=18, bold=True, color=TEAL)
        add_text(s, Inches(4.3), y + Inches(0.18), Inches(8.1), Inches(0.35), sub, size=14, bold=True, color=INK)
        add_bullets(s, Inches(0.9), y + Inches(0.6), Inches(11.5), Inches(0.95), items, size=13)
    footer(s, 14)

    # 15 对那段话的改写
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, WHITE)
    header(s, "把那段话改写成可用的版本", "保留心安，拿掉斩根。落地生根，不必切断旧根")
    add_round(s, Inches(0.5), Inches(1.2), Inches(6.1), Inches(5.5), INK)
    add_text(s, Inches(0.8), Inches(1.45), Inches(5.5), Inches(0.4), "原句（诀别）", size=14, bold=True, color=TEAL_LIGHT)
    add_text(
        s,
        Inches(0.8),
        Inches(2.0),
        Inches(5.5),
        Inches(4.3),
        "不要回来\n不要想念我们\n不要回头\n不要写信\n不要向乡愁屈服\n不要落叶归根\n而要落地生根\n心安之处即是家",
        size=18,
        color=WHITE,
    )
    add_round(s, Inches(6.8), Inches(1.2), Inches(6.05), Inches(5.5), MINT)
    add_text(s, Inches(7.1), Inches(1.45), Inches(5.5), Inches(0.4), "改写（战略）", size=14, bold=True, color=TEAL)
    add_text(
        s,
        Inches(7.1),
        Inches(2.0),
        Inches(5.5),
        Inches(4.3),
        "可以走，但要带着根走\n可以想家，但不让想家做决定\n可以回头看，但每周只看一次\n可以写信，把思念变成连接\n乡愁不是软肋，失控的乡愁才是\n不必归根，也不必斩根\n先在能创造价值的地方生根\n心安处即吾乡——心安要自己养",
        size=16,
        color=INK,
    )
    footer(s, 15)

    # 16 结语
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, INK)
    add_rect(s, 0, 0, Inches(0.18), SH, TEAL_LIGHT)
    add_text(s, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.35), "结语", size=14, color=TEAL_LIGHT)
    add_text(s, Inches(0.8), Inches(1.65), Inches(11.8), Inches(1.6), "别回头，向前走。\n不归根，也可以不斩根。\n心安处，即吾乡。", size=28, bold=True, color=WHITE)
    add_rect(s, Inches(0.8), Inches(3.5), Inches(2.2), Inches(0.06), TEAL_LIGHT)
    add_text(
        s,
        Inches(0.8),
        Inches(3.8),
        Inches(11.5),
        Inches(2.2),
        "2026 年的中国，仍适合一种发展：把人工智能接到真实的产业、空间与组织里。\n这段机会对转化者比对逃亡者更友好。美国仍是前沿高地，但城门在收；\n没有具体机会就「果断离开」，是把乡愁问题做成了错误的战略。\n先在能创造价值的地方把心安养出来。家，是养出来的，不是逃出来的。",
        size=16,
        color=SOFT,
    )
    add_text(s, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.4), "数据截至 2026 年 8 月 · 详见 Word 附录与 Excel 决策表", size=12, color=TEAL_LIGHT)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"已生成：{output_path}")


if __name__ == "__main__":
    root = Path(__file__).resolve().parents[1]
    build_ppt(root / "deliverables" / "发展与去留_2026中美判断备忘录.pptx")
