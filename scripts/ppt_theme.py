# -*- coding: utf-8 -*-
"""复兴岛策划方案 PPT 共用视觉模板（深蓝 + 青绿 + 金）。"""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

NAVY = RGBColor(0x0B, 0x1F, 0x3A)
NAVY2 = RGBColor(0x12, 0x2A, 0x4D)
TEAL = RGBColor(0x0A, 0x6E, 0x6A)
TEAL2 = RGBColor(0x1A, 0x9B, 0x8E)
GOLD = RGBColor(0xC4, 0x8A, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xE6, 0xF0, 0xEE)
GREY = RGBColor(0x9A, 0xA7, 0xBD)
CARD = RGBColor(0x15, 0x2E, 0x45)
ROW_ALT = RGBColor(0x1A, 0x36, 0x50)

FONT = "Microsoft YaHei"


def new_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def set_cjk(run, name=FONT):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def add_slide(prs, bg=NAVY):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid()
    r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def rect(s, x, y, w, h, color, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if color is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    return sp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space_after=4, line_spacing=1.05):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for t, size, color, bold in para:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            set_cjk(r)
    return tb


def header(s, kicker, title, idx, total, sw):
    rect(s, 0, 0, sw, Inches(1.12), NAVY2)
    rect(s, 0, Inches(1.12), sw, Pt(3), TEAL)
    rect(s, Inches(0.55), Inches(0.28), Pt(6), Inches(0.55), GOLD)
    text(s, Inches(0.75), Inches(0.16), Inches(10.2), Inches(0.85),
         [[(kicker, 11, TEAL2, True)], [(title, 21, WHITE, True)]], space_after=2)
    text(s, Inches(11.3), Inches(0.35), Inches(1.6), Inches(0.5),
         [[(f"{idx:02d}/{total:02d}", 14, GOLD, True)]], align=PP_ALIGN.RIGHT)


def footer(s, left_text, right_text="建议日期 2026-09-12"):
    text(s, Inches(0.7), Inches(7.08), Inches(8.8), Inches(0.3),
         [[(left_text, 9, GREY, False)]])
    text(s, Inches(9.5), Inches(7.08), Inches(3.2), Inches(0.3),
         [[(right_text, 9, GREY, False)]], align=PP_ALIGN.RIGHT)


def card(s, x, y, w, h, title, lines, accent=TEAL, tsize=14, bsize=11):
    rect(s, x, y, w, h, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, y, w, Pt(4), accent, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    runs = [[(title, tsize, WHITE, True)]]
    for ln in lines:
        runs.append([("· " + ln, bsize, LIGHT, False)])
    text(s, x + Inches(0.16), y + Inches(0.14), w - Inches(0.32), h - Inches(0.28),
         runs, space_after=2, line_spacing=1.02)


def add_table(s, left, top, width, height, headers, rows, col_widths=None, font_size=11):
    ncol = len(headers)
    nrow = 1 + len(rows)
    table_shape = s.shapes.add_table(nrow, ncol, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(font_size)
                r.font.bold = True
                r.font.color.rgb = WHITE
                set_cjk(r)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD if i % 2 == 0 else ROW_ALT
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if j else PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.size = Pt(font_size - 1)
                    r.font.color.rgb = LIGHT
                    set_cjk(r)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return table_shape
