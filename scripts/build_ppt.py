# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from content import AGENT_JOIN, DISCLAIMER, EVENT, FORWARD_SHORT, GUESTS, SERIES_NOTES, WATCH_FLOW, WHY_WATCH
from ppt_utils import (
    BLUE,
    CARD,
    CYAN,
    GOLD,
    MUTED,
    NAVY2,
    PINK,
    SLIDE_H,
    SLIDE_W,
    WHITE,
    add_paras,
    add_round,
    add_text,
    footer,
    header,
    paint_bg,
)

OUT = Path(__file__).resolve().parents[1] / "exports" / "飞书DemoDay4_观摩速览.pptx"


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    return s


def cover(prs):
    s = blank(prs)
    add_round(s, Inches(0.5), Inches(0.45), Inches(3.4), Inches(0.42), BLUE, 0.3)
    add_text(
        s,
        Inches(0.5),
        Inches(0.5),
        Inches(3.4),
        Inches(0.32),
        "AI BUILDER 系列  ·  第 4 场",
        size=12,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_text(s, Inches(0.5), Inches(1.15), Inches(12.2), Inches(0.4), "飞书 Demo Day #4", size=20, color=CYAN, bold=True)
    add_text(s, Inches(0.5), Inches(1.55), Inches(12.2), Inches(1.05), "让 AI 接住你的日常小麻烦", size=40, color=WHITE, bold=True)
    add_text(
        s,
        Inches(0.5),
        Inches(2.7),
        Inches(12),
        Inches(0.45),
        "Builder 现场演示  ·  张咋啦 / 向阳乔木  ·  新能力：Agent 入会",
        size=18,
        color=GOLD,
    )
    cards = [
        ("时间", "2026-07-02（周四）\n11:00–12:30 GMT+8"),
        ("形式", "线上直播\n飞书日历分享入会"),
        ("组织", "Ni Dan\n官方分享日历"),
        ("状态", "日历显示已结束\n先点链接看回放/加场"),
    ]
    for i, (k, v) in enumerate(cards):
        x = Inches(0.5 + i * 3.15)
        add_round(s, x, Inches(3.4), Inches(3.0), Inches(1.9), CARD)
        add_text(s, x + Inches(0.18), Inches(3.52), Inches(2.64), Inches(0.28), k, size=12, color=CYAN, bold=True)
        add_text(s, x + Inches(0.18), Inches(3.88), Inches(2.64), Inches(1.2), v, size=16, color=WHITE, bold=True)
    add_text(s, Inches(0.5), Inches(5.5), Inches(12.2), Inches(0.45), EVENT["share_url"], size=14, color=MUTED)
    add_text(s, Inches(0.5), Inches(6.05), Inches(12.2), Inches(0.85), DISCLAIMER, size=12, color=MUTED)
    footer(s, 1)


def why(prs):
    s = blank(prs)
    header(s, "为什么现在转发", "这场值得占坑的四个理由")
    labels = ["现场感", "可抄作业", "嘉宾密度", "新能力"]
    for i, (label, text) in enumerate(zip(labels, WHY_WATCH[:4])):
        y = Inches(1.28 + i * 1.18)
        add_round(s, Inches(0.5), y, Inches(12.3), Inches(1.08), CARD)
        add_round(s, Inches(0.7), y + Inches(0.28), Inches(1.9), Inches(0.5), BLUE, 0.3)
        add_text(
            s,
            Inches(0.7),
            y + Inches(0.34),
            Inches(1.9),
            Inches(0.38),
            label,
            size=14,
            color=WHITE,
            align=PP_ALIGN.CENTER,
            bold=True,
        )
        add_text(s, Inches(2.85), y + Inches(0.28), Inches(9.7), Inches(0.55), text, size=18, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 2)


def facts(prs):
    s = blank(prs)
    header(s, "已核对事实", "来自官方分享日历，不是猜测")
    rows = [
        ("标题", EVENT["title_share"]),
        ("开始", "2026 年 7 月 2 日 11:00（北京时间）"),
        ("结束", "2026 年 7 月 2 日 12:30（90 分钟）"),
        ("组织者", EVENT["organizer"]),
        ("入会方式", "打开分享链接 → 加入日历 / 进入会议"),
        ("日历状态", EVENT["status_zh"]),
        ("分享令牌", EVENT["token"]),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(1.25 + i * 0.72)
        add_round(s, Inches(0.5), y, Inches(2.4), Inches(0.62), NAVY2)
        add_round(s, Inches(3.05), y, Inches(9.75), Inches(0.62), CARD)
        add_text(s, Inches(0.62), y + Inches(0.12), Inches(2.2), Inches(0.4), k, size=16, color=CYAN, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(3.2), y + Inches(0.12), Inches(9.4), Inches(0.4), v, size=16, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 3)


def series(prs):
    s = blank(prs)
    header(s, "系列是什么", "AI Builder Demo Day，不是发布会")
    add_round(s, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.4), CARD)
    add_text(s, Inches(0.75), Inches(1.5), Inches(5.5), Inches(0.4), "可以确定的事", size=18, color=GOLD, bold=True)
    add_paras(s, Inches(0.75), Inches(2.05), Inches(5.5), Inches(4.4), SERIES_NOTES, size=15)
    add_round(s, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.4), CARD)
    add_text(s, Inches(7.05), Inches(1.5), Inches(5.5), Inches(0.4), "观摩时不要误判", size=18, color=PINK, bold=True)
    add_paras(
        s,
        Inches(7.05),
        Inches(2.05),
        Inches(5.5),
        Inches(4.4),
        [
            "这不是飞书年度发布会，是社区 Builder 路演。",
            "标题里的「2」更像加场，不要理解成 Demo Day 第 2 季。",
            "嘉宾名单来自转发预告，现场以主持人为准。",
            "Agent 入会是能力演示，不等于明天就能进你的正式会。",
        ],
        size=15,
    )
    footer(s, 4)


def theme(prs):
    s = blank(prs)
    header(s, "主题怎么用", "日常小麻烦 > 炫技大工程")
    items = [
        ("会前", "纪要、议程、材料还在十个群里飞"),
        ("会中", "有人没来、有人抢话、没人记待办"),
        ("会后", "纪要隔两小时才出，待办对不齐人"),
        ("重复", "每周同样的站会、同样的周报、同样的催"),
    ]
    for i, (k, v) in enumerate(items):
        x = Inches(0.5 + (i % 2) * 6.4)
        y = Inches(1.35 + (i // 2) * 2.5)
        add_round(s, x, y, Inches(6.1), Inches(2.25), CARD)
        add_text(s, x + Inches(0.3), y + Inches(0.3), Inches(5.5), Inches(0.4), k, size=16, color=CYAN, bold=True)
        add_text(s, x + Inches(0.3), y + Inches(0.85), Inches(5.5), Inches(1.05), v, size=22, color=WHITE, bold=True)
    footer(s, 5)


def guests(prs):
    s = blank(prs)
    header(s, "预告嘉宾", "对照看，不追星")
    for i, g in enumerate(GUESTS):
        x = Inches(0.5 + i * 6.4)
        add_round(s, x, Inches(1.3), Inches(6.1), Inches(5.4), CARD)
        add_text(s, x + Inches(0.3), Inches(1.5), Inches(5.5), Inches(0.45), g["name"], size=26, color=WHITE, bold=True)
        add_text(s, x + Inches(0.3), Inches(2.0), Inches(5.5), Inches(0.3), f"{g['aka']}  ·  {g['role']}", size=13, color=GOLD)
        add_text(s, x + Inches(0.3), Inches(2.45), Inches(5.5), Inches(1.35), g["why"], size=14, color=MUTED)
        add_text(s, x + Inches(0.3), Inches(3.9), Inches(5.5), Inches(0.3), "现场盯这三句", size=14, color=CYAN, bold=True)
        add_paras(s, x + Inches(0.3), Inches(4.25), Inches(5.5), Inches(2.1), g["watch"], size=14)
    footer(s, 6)


def agent(prs):
    s = blank(prs)
    header(s, "本场最大变量", "Agent 入会：AI 从会后总结变成会中同事")
    add_round(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.35), CARD)
    add_text(s, Inches(0.75), Inches(1.5), Inches(11.8), Inches(1.0), AGENT_JOIN["what"], size=18, color=WHITE)
    for i, (title, items, color) in enumerate(
        [
            ("它能做什么", AGENT_JOIN["can"], CYAN),
            ("先别指望", AGENT_JOIN["cannot"], PINK),
        ]
    ):
        x = Inches(0.5 + i * 6.4)
        add_round(s, x, Inches(2.85), Inches(6.1), Inches(3.8), CARD)
        add_text(s, x + Inches(0.3), Inches(3.05), Inches(5.5), Inches(0.4), title, size=18, color=color, bold=True)
        add_paras(s, x + Inches(0.3), Inches(3.55), Inches(5.5), Inches(2.85), items, size=15)
    footer(s, 7)


def agent_q(prs):
    s = blank(prs)
    header(s, "会中只问这 6 个问题", "把 Demo 听成可落地的能力清单")
    for i, q in enumerate(AGENT_JOIN["watch_questions"]):
        y = Inches(1.28 + i * 0.88)
        add_round(s, Inches(0.5), y, Inches(12.3), Inches(0.78), CARD)
        add_round(s, Inches(0.7), y + Inches(0.16), Inches(0.52), Inches(0.46), BLUE, 0.2)
        add_text(
            s,
            Inches(0.7),
            y + Inches(0.2),
            Inches(0.52),
            Inches(0.38),
            f"{i+1:02d}",
            size=13,
            color=WHITE,
            align=PP_ALIGN.CENTER,
            bold=True,
        )
        add_text(s, Inches(1.45), y + Inches(0.18), Inches(11.1), Inches(0.45), q, size=18, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 8)


def flow(prs):
    s = blank(prs)
    header(s, "观摩节奏", "会前占坑 · 会中抓步骤 · 会后 24 小时复现")
    for i, (title, items) in enumerate(WATCH_FLOW.items()):
        x = Inches(0.45 + i * 4.25)
        add_round(s, x, Inches(1.3), Inches(4.05), Inches(5.4), CARD)
        add_text(s, x + Inches(0.25), Inches(1.5), Inches(3.55), Inches(0.45), title, size=22, color=GOLD, bold=True)
        add_paras(s, x + Inches(0.25), Inches(2.15), Inches(3.55), Inches(4.2), items, size=14)
    footer(s, 9)


def soldout(prs):
    s = blank(prs)
    header(s, "上一场的教训", "先加会议，再讨论看不看")
    steps = [
        ("1", "立刻点开分享链接", "加入日历，不要等开场前五分钟"),
        ("2", "确认入会权限", "外部账号、公司租户、是否要审批"),
        ("3", "设两个提醒", "开场前 30 分钟、前 5 分钟"),
        ("4", "约满就换入口", "B 站搜官方回放，或等加场"),
    ]
    for i, (n, t, d) in enumerate(steps):
        y = Inches(1.3 + i * 1.2)
        add_round(s, Inches(0.5), y, Inches(12.3), Inches(1.08), CARD)
        add_round(s, Inches(0.7), y + Inches(0.24), Inches(0.6), Inches(0.6), BLUE, 0.2)
        add_text(s, Inches(0.7), y + Inches(0.32), Inches(0.6), Inches(0.45), n, size=18, color=WHITE, align=PP_ALIGN.CENTER, bold=True)
        add_text(s, Inches(1.55), y + Inches(0.18), Inches(10.9), Inches(0.4), t, size=20, color=WHITE, bold=True)
        add_text(s, Inches(1.55), y + Inches(0.58), Inches(10.9), Inches(0.35), d, size=15, color=MUTED)
    footer(s, 10)


def copy_slide(prs):
    s = blank(prs)
    header(s, "可直接发群", "短口径，保留预约链接")
    add_round(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.4), CARD)
    add_text(s, Inches(0.8), Inches(1.55), Inches(11.7), Inches(4.9), FORWARD_SHORT, size=20, color=WHITE)
    footer(s, 11)


def actions(prs):
    s = blank(prs)
    header(s, "一页行动", "看完这场，至少带走一件可做的事")
    items = [
        ("占坑", "把日历链接发给要一起看的人，而不是开场再喊"),
        ("选题", "各自写一个「最烦的会」：周会 / 客户会 / 评审会"),
        ("对照", "现场只记：麻烦、人、Agent、权限、出口文档"),
        ("复现", "会后用内部测试会走一遍 Agent 入会，不上正式会"),
        ("沉淀", "把可抄步骤写成飞书文档，变成下一次 brief"),
        ("补课", "没看成直播就去 B 站搜 AI Builder Demo Day"),
    ]
    for i, (k, v) in enumerate(items):
        x = Inches(0.5 + (i % 3) * 4.2)
        y = Inches(1.3 + (i // 3) * 2.55)
        add_round(s, x, y, Inches(4.0), Inches(2.35), CARD)
        add_text(s, x + Inches(0.25), y + Inches(0.25), Inches(3.5), Inches(0.4), k, size=20, color=GOLD, bold=True)
        add_text(s, x + Inches(0.25), y + Inches(0.8), Inches(3.5), Inches(1.25), v, size=16, color=WHITE)
    footer(s, 12)


def appendix(prs):
    s = blank(prs)
    header(s, "附录 · 会后自己试", "公开能力：飞书 CLI 的会中机器人")
    add_round(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.5), CARD)
    add_text(s, Inches(0.75), Inches(1.5), Inches(11.8), Inches(0.35), "这是公开技能说明，不是本场官方议程。", size=14, color=GOLD)
    add_text(s, Inches(0.75), Inches(1.95), Inches(11.8), Inches(0.55), AGENT_JOIN["join_cmd"], size=18, color=WHITE, bold=True)
    blocks = [
        ("入会", "9 位会议号 + 机器人身份。会议必须进行中。"),
        ("事件", "用返回的 meeting.id 拉进出 / 发言 / 共享。"),
        ("发言", "会中文本或表情，别拿来刷屏。"),
        ("离会", "看完就退，避免机器人一直挂在会里。"),
    ]
    for i, (k, v) in enumerate(blocks):
        x = Inches(0.5 + i * 3.15)
        add_round(s, x, Inches(3.05), Inches(3.0), Inches(2.85), CARD)
        add_text(s, x + Inches(0.2), Inches(3.25), Inches(2.6), Inches(0.4), k, size=18, color=CYAN, bold=True)
        add_text(s, x + Inches(0.2), Inches(3.8), Inches(2.6), Inches(1.8), v, size=15, color=WHITE)
    footer(s, 13)


def end(prs):
    s = blank(prs)
    add_text(s, Inches(0.5), Inches(2.0), Inches(12.2), Inches(0.4), "先加会议。", size=20, color=CYAN, bold=True)
    add_text(s, Inches(0.5), Inches(2.5), Inches(12.2), Inches(0.9), "上一场爆满，后面约不上。", size=36, color=WHITE, bold=True)
    add_text(s, Inches(0.5), Inches(3.6), Inches(12.2), Inches(0.6), EVENT["share_url"], size=16, color=GOLD)
    add_text(
        s,
        Inches(0.5),
        Inches(4.5),
        Inches(12.2),
        Inches(1.3),
        "材料用途：内部转发、会前预习、会中记录、会后复盘。\n不是飞书官方议程，现场以主持人和日历为准。",
        size=16,
        color=MUTED,
    )
    footer(s, 14)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    cover(prs)
    why(prs)
    facts(prs)
    series(prs)
    theme(prs)
    guests(prs)
    agent(prs)
    agent_q(prs)
    flow(prs)
    soldout(prs)
    copy_slide(prs)
    actions(prs)
    appendix(prs)
    end(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"wrote {OUT} slides={len(prs.slides)}")
    return OUT


if __name__ == "__main__":
    build()
