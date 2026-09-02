"""
冠松 01# · 报价与论证发给对方的 6 页 PPT

重新生成：python3 scripts/build_quote_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pathlib import Path

NAVY = RGBColor(0x1B, 0x35, 0x55)
GOLD = RGBColor(0xB7, 0x86, 0x2E)
INK = RGBColor(0x2A, 0x2D, 0x34)
CLOUD = RGBColor(0xF0, 0xF1, 0xF3)
GREY = RGBColor(0x6B, 0x73, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0x8B, 0x2E, 0x2E)
MOSS = RGBColor(0x2E, 0x5A, 0x3C)

CN_FONT = "WenQuanYi Micro Hei"
EN_FONT = "Georgia"
SLIDES = []


def set_run(run, text, *, size=14, bold=False, color=INK, italic=False):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = EN_FONT
    rPr = run._r.get_or_add_rPr()
    for tag in ("ea", "cs"):
        existing = rPr.find(qn(f"a:{tag}"))
        if existing is not None:
            rPr.remove(existing)
        rPr.append(rPr.makeelement(qn(f"a:{tag}"), {"typeface": CN_FONT}))


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    box.fill.background()
    box.line.fill.background()
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        set_run(run, ln, size=size, bold=bold, color=color, italic=italic)
    return box


def add_rect(slide, x, y, w, h, *, fill=NAVY):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    shp.text_frame.text = ""
    return shp


def add_round(slide, x, y, w, h, text="", *, fill=NAVY, color=WHITE,
              size=12, bold=True, align=PP_ALIGN.CENTER):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p.add_run(), text, size=size, bold=bold, color=color)
    return shp


def add_table(slide, x, y, w, h, header, rows, *, header_size=12, body_size=12,
              col_widths=None):
    ts = slide.shapes.add_table(len(rows) + 1, len(header), x, y, w, h)
    table = ts.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw
    for j, htxt in enumerate(header):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.text = ""
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p.add_run(), htxt, size=header_size, bold=True, color=WHITE)
    for i, row in enumerate(rows, start=1):
        bg = WHITE if i % 2 else CLOUD
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            cell.text = ""
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            set_run(p.add_run(), str(val), size=body_size, color=INK)
    return ts


def add_chrome(slide, prs, *, page_no, label, title, subtitle=""):
    sw, sh = prs.slide_width, prs.slide_height
    add_rect(slide, 0, 0, sw, Emu(380000), fill=NAVY)
    add_round(slide, Inches(0.45), Inches(0.18), Inches(2.4), Inches(0.36),
              label, fill=GOLD, color=NAVY, size=11)
    add_text(slide, Inches(3.05), Inches(0.10), Inches(9.6), Inches(0.52),
             title, size=22, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(3.05), Inches(0.52), Inches(9.6), Inches(0.28),
                 subtitle, size=12, color=CLOUD, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 0, sh - Emu(80000), sw, Emu(80000), fill=GOLD)
    add_text(slide, Inches(0.45), Inches(7.05), Inches(9.2), Inches(0.28),
             "冠松 01# 楼 · 报价函配套 · 保密", size=10, color=GREY)
    add_text(slide, Inches(11.4), Inches(7.05), Inches(1.5), Inches(0.28),
             f"{page_no} / 0", size=10, color=GREY, align=PP_ALIGN.RIGHT)


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    SLIDES.append(s)
    return s


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    sw, sh = prs.slide_width, prs.slide_height

    s = new_slide(prs)
    add_rect(s, 0, 0, sw, sh, fill=NAVY)
    add_rect(s, 0, Inches(5.05), sw, Emu(28000), fill=GOLD)
    add_text(s, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.4),
             "01# 研发楼 · 2026-09-03 报价",
             size=16, bold=True, color=GOLD, italic=True)
    add_text(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(1.3),
             "政府公关为主，招商为辅",
             size=36, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.2), Inches(11.7), Inches(1.2),
             "配合把 C6 往公寓 / 酒店的路径问清楚\n不保证能批 · 补差价随时同步 · 招商同步不放",
             size=18, color=CLOUD)
    add_round(s, Inches(0.8), Inches(5.5), Inches(5.6), Inches(0.48),
              "启动金 15 万覆盖头 90 天", fill=GOLD, color=NAVY, size=16)
    add_text(s, Inches(6.6), Inches(5.5), Inches(5.8), Inches(0.48),
             "之后 8 万 / 月 · 散 8% / 整层 12%",
             size=16, color=CLOUD, anchor=MSO_ANCHOR.MIDDLE)

    s = new_slide(prs)
    add_chrome(s, prs, page_no=2, label="论证",
               title="C6 改公寓酒店：先拆开，再去问",
               subtitle="论证的不是「已经能改」，而是「有几扇门、过几道坎、哪一道可能补差价」。")
    add_table(
        s, Inches(0.4), Inches(1.25), Inches(12.5), Inches(4.35),
        ["形态", "预判", "90 天"],
        [
            ["整栋持牌酒店", "低 · 多半改性质、补差大", "只问规则和账，不做承诺"],
            ["整栋市场化公寓 / 商住", "很低 · 易被看成改住宅", "点明风险，不主攻"],
            ["部分楼层服务式公寓（整体持有）", "低到中 · 看更新窗口", "有窗口再加深，先装修不算"],
            ["人才公寓 / 专家公寓（科研配套）", "相对可争 · 仍不保证", "主攻：不改主导功能、不分割"],
            ["维持研发 + 小型配套", "高", "保底；招商随时可执行"],
        ],
        header_size=14, body_size=13,
        col_widths=[Inches(4.5), Inches(4.0), Inches(4.0)],
    )
    add_rect(s, Inches(0.4), Inches(5.75), Inches(12.5), Inches(0.9), fill=CLOUD)
    add_text(s, Inches(0.55), Inches(5.85), Inches(12.2), Inches(0.7),
             "静安「酒店、人才公寓」写在存量商务楼宇条款里。01# 是 C6 新建研发楼，该走园区更新这扇门，不能直接套。",
             size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    s = new_slide(prs)
    add_chrome(s, prs, page_no=3, label="方法",
               title="五步论证 · 用请示问，不用方案压",
               subtitle="书面批复不做保证。没有纪要和染色，才算没做完。")
    steps = [
        ("1", "读地证", "合同、图则、消防专篇。没有原文不去空谈。"),
        ("2", "选对门", "园区更新，不是商务楼宇更新。M0 到不了酒店。"),
        ("3", "四条请示", "入口 · 配套住宿 · 酒店是否改性质 · 与 9/30 关系。"),
        ("4", "消防单列", "土地松口 ≠ 能开业。同项目不同楼可以相反。"),
        ("5", "补差单独账", "要不要补、规则、周期。数字一出立刻同步。"),
    ]
    for i, (n, t, b) in enumerate(steps):
        y = Inches(1.22) + Inches(1.05) * i
        add_round(s, Inches(0.5), y, Inches(0.7), Inches(0.88), n,
                  fill=NAVY, color=WHITE, size=20)
        add_text(s, Inches(1.4), y, Inches(2.4), Inches(0.88), t,
                 size=18, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(3.9), y, Inches(8.8), Inches(0.88), b,
                 size=16, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    s = new_slide(prs)
    add_chrome(s, prs, page_no=4, label="边界",
               title="能配合，不保证；补差价随时说",
               subtitle="不设获批成功费，避免报价看起来像在卖批文。")
    left = [
        ("我们配合", MOSS,
         "写请示、备材料、陪同沟通、做纪要\n路径染色交给董事会表决\n招商备选库和渠道同步推进"),
        ("我们不保证", RED,
         "不保证能改酒店或市场化公寓\n不保证不补地价、不保证补多少\n不保证规划消防营业一定过"),
    ]
    for i, (t, c, b) in enumerate(left):
        x = Inches(0.45) + Inches(6.4) * i
        add_rect(s, x, Inches(1.25), Inches(6.1), Inches(3.15), fill=c)
        add_text(s, x, Inches(1.4), Inches(6.1), Inches(0.5), t,
                 size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.25), Inches(2.05), Inches(5.6), Inches(2.1),
                 b, size=16, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(0.45), Inches(4.6), Inches(12.4), Inches(2.0), fill=CLOUD)
    add_text(s, Inches(0.7), Inches(4.75), Inches(12.0), Inches(1.7),
             "补差价是土地使用权价款，主体是业主。我们问清规则和周期，数字一经窗口或估价机构提出，立即书面同步。\n"
             "减免、分期可协助准备口径，不代为承诺。估价、测绘、设计消防变更，不包在启动金里。\n"
             "9 月 30 日按已批研发楼进场。政策并行，不拿滞纳金赌博。",
             size=15, color=INK)

    s = new_slide(prs)
    add_chrome(s, prs, page_no=5, label="报价",
               title="两线做在一笔启动金里",
               subtitle="不另收公关项目费。头 90 天现金就是 15 万。招到了，15 万抵佣金。")
    add_table(
        s, Inches(0.4), Inches(1.22), Inches(12.5), Inches(3.85),
        ["项", "数", "说明"],
        [
            ["启动金", "15 万 / 90 天", "论证、请示、纪要、可决稿；同步招商备选与渠道"],
            ["月度费", "8 万 / 月", "第 91 天起，两线继续；当月有佣金则抵扣"],
            ["散户佣金", "首年租金 8%", "生效租赁才收；甲方自有未报备客户不收"],
            ["整层 / 大客户", "首年租金 12%", "9 月 2 日见面口径"],
            ["孵化器品牌", "另计", "不包在 15 万内"],
        ],
        header_size=14, body_size=13,
        col_widths=[Inches(2.6), Inches(3.2), Inches(6.7)],
    )
    add_rect(s, Inches(0.4), Inches(5.25), Inches(12.5), Inches(1.35), fill=NAVY)
    add_text(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.1),
             "招商不放弃：政策在问的同时，研发客户库照建，避免楼盖好了两边都空。无协议不代看。\n"
             "常规中介约年租金 16%；我方佣金大约一半量级。中介带来的单，我方只收一半比例。",
             size=15, color=WHITE)

    s = new_slide(prs)
    add_chrome(s, prs, page_no=6, label="请定",
               title="请确认这四项，我们按函推进",
               subtitle="纸质报价函有效期 15 个工作日。正式权利义务以盖章协议为准。")
    asks = [
        ("01", "费用", "15 万 / 90 天 · 之后 8 万/月 · 散 8% · 整层 12%。"),
        ("02", "边界", "配合推进，不保证获批；补差价随时沟通，业主承担并决策。"),
        ("03", "主攻", "配套住宿论证为主；整栋酒店只问；招商同步不放。"),
        ("04", "材料", "3 个工作日内土地证 + 出让合同用途全文；指定出面人。"),
    ]
    for i, (n, k, v) in enumerate(asks):
        y = Inches(1.25) + Inches(1.25) * i
        add_round(s, Inches(0.5), y, Inches(1.15), Inches(1.05), n,
                  fill=GOLD, color=NAVY, size=20)
        add_text(s, Inches(1.9), y, Inches(2.0), Inches(1.05), k,
                 size=22, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(4.05), y, Inches(8.6), Inches(1.05), v,
                 size=18, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    total = len(SLIDES)
    for sl in SLIDES:
        for shape in sl.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text and run.text.strip().endswith(" / 0"):
                        run.text = run.text.replace(" / 0", f" / {total}")

    out = Path(__file__).resolve().parent.parent / "docs" / "advisory" / \
        "deck" / "冠松01楼-报价与论证.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"✓ Deck written: {out}  ({total} slides)")
    return out


if __name__ == "__main__":
    build()
