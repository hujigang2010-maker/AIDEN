# -*- coding: utf-8 -*-
"""
生成《上海市创业孵化示范基地 · 2026年工作报告》PPT。

结构对标 2023 年评审 PPT，数据与表述依据 2026 年度工作报告（覆盖 2025 年工作）。
图片复用自 2023 年 PPT/PDF 中的园区活动与载体实景素材。

运行：python3 scripts/build_2026_work_report_ppt.py
输出：dist/复旦科技园创业孵化基地2026年工作报告.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
IMG = ROOT / "assets" / "images"
DIST = ROOT / "dist"
DIST.mkdir(parents=True, exist_ok=True)
OUT = DIST / "复旦科技园创业孵化基地2026年工作报告.pptx"

# 复旦科技园惯用深蓝体系（对标原 PPT）
NAVY = RGBColor(0x0A, 0x2F, 0x6B)
NAVY_DK = RGBColor(0x07, 0x1F, 0x4A)
BLUE = RGBColor(0x1F, 0x4E, 0x8C)
LIGHT = RGBColor(0xF5, 0xF8, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1B, 0x2A, 0x3A)
MUTED = RGBColor(0x5B, 0x6B, 0x7C)
RED = RGBColor(0xC0, 0x39, 0x2B)
GOLD = RGBColor(0xC9, 0xA2, 0x4B)
CARD = RGBColor(0xE8, 0xEF, 0xF8)

FONT = "微软雅黑"
SW_IN, SH_IN = 13.333, 7.5

prs = Presentation()
prs.slide_width = Inches(SW_IN)
prs.slide_height = Inches(SH_IN)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _font(run, size, color, bold=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = FONT
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT)


def add_slide(bg=WHITE):
    slide = prs.slides.add_slide(BLANK)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg
    rect.line.fill.background()
    return slide


def rect(slide, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: (text, size, color, bold[, space_after])"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    for i, spec in enumerate(lines):
        t, size, color, bold = spec[0], spec[1], spec[2], spec[3]
        space_after = spec[4] if len(spec) > 4 else 4
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = t
        _font(run, size, color, bold)
    return tb


def picture(slide, name, x, y, w, h=None):
    path = IMG / name
    if not path.exists():
        # try alternate extensions
        for ext in (".jpeg", ".jpg", ".png", ".webp"):
            alt = IMG / f"{Path(name).stem}{ext}"
            if alt.exists():
                path = alt
                break
    if not path.exists():
        return None
    if h is None:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def header(slide, title, page_no=None):
    rect(slide, 0, 0, SW_IN, 0.95, NAVY)
    # 优先透明红标，缺失时回退
    if not picture(slide, "logo_red_mark.png", 0.28, 0.2, 0.5, 0.5):
        picture(slide, "logo_red_mark.jpeg", 0.28, 0.2, 0.5, 0.5)
    text(
        slide,
        0.9,
        0.18,
        4.5,
        0.6,
        [("复旦大学国家大学科技园", 12, WHITE, False)],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    text(
        slide,
        5.5,
        0.2,
        7.3,
        0.55,
        [(title, 22, WHITE, True)],
        align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    rect(slide, 0, 7.25, SW_IN, 0.25, NAVY)
    if page_no:
        text(
            slide,
            12.2,
            7.22,
            0.9,
            0.28,
            [(str(page_no), 10, WHITE, False)],
            align=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE,
        )


def card(slide, x, y, w, h, title, body_lines, accent=BLUE):
    rect(slide, x, y, w, h, CARD)
    rect(slide, x, y, 0.1, h, accent)
    text(slide, x + 0.25, y + 0.15, w - 0.4, 0.4, [(title, 14, NAVY, True)])
    lines = [(b, 11, DARK, False, 3) for b in body_lines]
    text(slide, x + 0.25, y + 0.55, w - 0.4, h - 0.7, lines)


# ============================================================================
# 1 封面
# ============================================================================
s = add_slide(WHITE)
picture(s, "cover_building.jpeg", 6.2, 0, 7.2, 7.5)
rect(s, 0, 0, 7.4, 7.5, WHITE)
rect(s, 6.9, 0, 0.35, 7.5, NAVY)
rect(s, 0, 6.85, SW_IN, 0.65, NAVY)
# 品牌区：红标 + 校名
if not picture(s, "logo_red_mark.png", 0.7, 0.65, 0.72, 0.72):
    picture(s, "logo_red_mark.jpeg", 0.7, 0.65, 0.72, 0.72)
text(s, 1.6, 0.65, 4.5, 0.8, [
    ("复旦大学", 20, NAVY, True, 0),
    ("国家大学科技园", 14, MUTED, False),
])
text(s, 0.7, 2.4, 6, 1.6, [
    ("2026年工作报告", 40, NAVY, True, 10),
    ("上海市创业孵化示范基地", 20, BLUE, False),
])
text(s, 0.7, 5.6, 5, 0.5, [("复旦科技园创业孵化基地", 14, MUTED, False)])
text(s, 0.7, 7.0, 5, 0.35, [("2026.06", 14, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)

# ============================================================================
# 2 目录
# ============================================================================
s = add_slide(LIGHT)
header(s, "目录", 2)
items = [
    ("01", "创业孵化功能", "载体定位 · 服务体系 · 导师与团队"),
    ("02", "创业孵化成效", "孵化实效 · 带动就业 · 政策融资"),
    ("03", "人社部门合作紧密度", "政策落实 · 联合活动 · 宣传协同"),
    ("04", "特色工作", "品牌赛事 · 成果转化 · 开放协同"),
]
for i, (no, title, sub) in enumerate(items):
    y = 1.4 + i * 1.3
    rect(s, 1.2, y, 10.8, 1.1, WHITE)
    rect(s, 1.2, y, 1.3, 1.1, NAVY)
    text(s, 1.2, y, 1.3, 1.1, [(no, 28, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 2.8, y + 0.2, 8.5, 0.45, [(title, 22, NAVY, True)])
    text(s, 2.8, y + 0.6, 8.5, 0.35, [(sub, 13, MUTED, False)])

# ============================================================================
# 3 载体基本情况
# ============================================================================
s = add_slide(WHITE)
header(s, "一、复旦科技园创业孵化功能", 3)
text(s, 0.55, 1.2, 7.0, 3.6, [
    ("复旦大学国家大学科技园是科技部、教育部联合认定的首批国家大学科技园，作为大学原创力向社会延伸的重要载体。", 13, DARK, False, 8),
    ("基地依托复旦科技园建设运营，秉持“转化科技、服务社会、汇聚智慧、共创未来”理念，构建“成果发现－概念验证－创业孵化－产业培育”全链条服务，打造高校科技成果转化、创新创业生态和区域发展新动能平台。", 13, DARK, False, 8),
    ("先后获评国家高新技术创业服务中心、上海市创业孵化示范基地、上海市海聚英才创新创业示范基地等资质。", 13, DARK, False, 8),
])
# KPI 条
kpis = [
    ("地址", "杨浦区国泰路11号"),
    ("场地面积", "3.4万㎡"),
    ("孵化面积", "2.56万㎡"),
    ("在孵实体", "159家"),
]
for i, (k, v) in enumerate(kpis):
    x = 0.55 + i * 2.95
    rect(s, x, 5.0, 2.75, 1.7, CARD)
    text(s, x + 0.15, 5.2, 2.45, 0.4, [(k, 12, MUTED, False)])
    text(s, x + 0.15, 5.65, 2.45, 0.7, [(v, 22, RED if i else NAVY, True)])
picture(s, "campus_building.jpeg", 7.9, 1.25, 4.8, 3.5)

# ============================================================================
# 4 四大功能
# ============================================================================
s = add_slide(WHITE)
header(s, "四大孵化功能板块", 4)
funcs = [
    ("01 创业苗圃", "聚焦师生创业孵化，配合复旦大学开展创新创业教育与实践。"),
    ("02 复翼互联众创空间", "搭建全生命周期孵化体系，联动学院、校友会输送项目与资源。"),
    ("03 技术转移服务平台", "通过技术转移中心、校企联合实验室对接科研学科，提供技术赋能。"),
    ("04 海外人才科创基地", "依托海聚英才等载体，吸引海外创新创业人才集聚发展。"),
]
for i, (t, b) in enumerate(funcs):
    x = 0.55 + (i % 2) * 6.3
    y = 1.35 + (i // 2) * 2.7
    rect(s, x, y, 6.0, 2.4, CARD)
    rect(s, x, y, 6.0, 0.12, NAVY)
    text(s, x + 0.35, y + 0.4, 5.3, 0.5, [(t, 18, NAVY, True)])
    text(s, x + 0.35, y + 1.1, 5.3, 1.0, [(b, 14, DARK, False)])

# ============================================================================
# 5 服务体系
# ============================================================================
s = add_slide(WHITE)
header(s, "创业孵化功能服务体系", 5)
text(s, 0.55, 1.2, 12, 0.4, [("建立规范化、体系化的日常管理工作体系", 16, NAVY, True)])
services = ["企业日常走访", "基础企业服务", "投融资与市场对接", "创新创业活动", "技术转移服务", "政策辅导申报"]
for i, name in enumerate(services):
    x = 0.55 + (i % 3) * 4.15
    y = 1.85 + (i // 3) * 1.5
    rect(s, x, y, 3.95, 1.25, CARD)
    text(s, x, y, 3.95, 1.25, [(name, 16, NAVY, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.55, 5.2, 12, 1.5, [
    ("管理制度覆盖：大学生创业基地管理办法、科创企业管理办法、投资管理办法、创业孵化管理办法、入孵与退出办法、促进科技成果转化实施办法及操作细则等。", 13, DARK, False, 6),
    ("2025年新增创业导师16人；签约第三方服务机构41家，全年服务超2000家次；复煜概念验证中心投入运营。", 13, DARK, False),
])

# ============================================================================
# 6 管理与导师
# ============================================================================
s = add_slide(WHITE)
header(s, "创业导师与基地管理团队", 6)
card(s, 0.55, 1.3, 6.0, 2.5, "创业导师队伍", [
    "高校教授、企业家、投资人、优秀校友多元构成",
    "开展创业讲堂、校友沙龙、一对一辅导等活动",
    "2025年新增创业导师16人，持续提升辅导深度",
])
card(s, 6.8, 1.3, 5.9, 2.5, "基地管理团队", [
    "具备科技园与孵化器长期运营经验",
    "对接人社创业指导、政策辅导与投融资服务",
    "累计服务孵化企业规模持续扩大",
], accent=GOLD)
picture(s, "mentor1.jpeg", 0.55, 4.1, 2.8, 2.7)
picture(s, "mentor2.jpeg", 3.55, 4.1, 2.8, 2.7)
picture(s, "activity1.jpeg", 6.8, 4.1, 5.9, 2.7)

# ============================================================================
# 7 成效总览
# ============================================================================
s = add_slide(WHITE)
header(s, "二、创业孵化成效", 7)
metrics = [
    ("新增入孵团队", "26家", "当年孵化成功率100%"),
    ("新增初创组织", "23家", "年末存活率100%"),
    ("在孵创业实体", "159家", "2025年末存续"),
    ("在岗人数", "324人", "同比+57.28%"),
    ("新增就业", "118人", "创业带动就业"),
    ("双创活动", "47场", "与人社合作3场"),
    ("投融资支持", "9700余万", "帮扶创业组织融资"),
    ("人社政策受益", "9家", "担保贷款2家"),
]
for i, (k, v, sub) in enumerate(metrics):
    x = 0.45 + (i % 4) * 3.2
    y = 1.35 + (i // 4) * 2.7
    rect(s, x, y, 3.0, 2.4, CARD)
    text(s, x + 0.2, y + 0.3, 2.6, 0.4, [(k, 13, MUTED, False)])
    text(s, x + 0.2, y + 0.85, 2.6, 0.7, [(v, 28, RED, True)])
    text(s, x + 0.2, y + 1.7, 2.6, 0.4, [(sub, 12, DARK, False)])

# ============================================================================
# 8 成效细分
# ============================================================================
s = add_slide(WHITE)
header(s, "创业孵化成效 · 重点指标", 8)
blocks = [
    ("创业组织培育", ["新增入孵创业团队26家、初创组织23家", "年末在孵创业实体159家", "优秀创业组织9家"]),
    ("选树创业典型", ["眸深智能获全国颠覆性技术创新大赛最高奖", "本科生团队获中国国际大学生创新大赛金奖", "音书科技获天使基金优秀雏鹰企业"]),
    ("创业氛围营造", ["双创活动47场，培训/路演/沙龙并举", "推荐赛事项目100余项", "宣传稿34篇，人社供稿采纳10篇"]),
    ("创业融资服务", ["投融资支持约9700余万元", "联动大创基金、科创母基金与前孵化资金", "合作银行与创投机构持续扩容"]),
]
for i, (t, lines) in enumerate(blocks):
    x = 0.5 + (i % 2) * 6.4
    y = 1.3 + (i // 2) * 2.8
    card(s, x, y, 6.1, 2.55, t, lines)

# ============================================================================
# 9 人社合作
# ============================================================================
s = add_slide(WHITE)
header(s, "三、与人社部门合作紧密度", 9)
text(s, 0.55, 1.2, 7.2, 2.8, [
    ("1. 政策宣传与对接落实", 16, NAVY, True, 6),
    ("围绕人社创业扶持政策开展宣传辅导，2025年推荐9家创业组织享受人社扶持政策，帮扶2家获得创业担保贷款。", 13, DARK, False, 10),
    ("2. 联合开展创业活动", 16, NAVY, True, 6),
    ("与人社部门合作开展创新创业活动3场；推荐组织和项目参加人社类活动/赛事12项；马兰花GYB培训班顺利举办。", 13, DARK, False, 10),
    ("3. 宣传协同", 16, NAVY, True, 6),
    ("向人社部门供稿11篇，其中10篇被“海纳百创”“乐业杨浦”等媒介采纳。", 13, DARK, False),
])
picture(s, "event_hr1.jpeg", 8.0, 1.25, 4.7, 2.7)
picture(s, "activity2.jpeg", 8.0, 4.15, 4.7, 2.7)

# ============================================================================
# 10 特色总览
# ============================================================================
s = add_slide(WHITE)
header(s, "四、特色工作", 10)
feats = [
    ("01", "助力高校创业教育", "课程实训、见习指导、马兰花培训与学分实践相结合"),
    ("02", "推动成果转移转化", "对接16个院系，入库成果21项，转化落地12项"),
    ("03", "就业促进与招聘", "校园招聘、见习实训与岗位对接同步推进"),
    ("04", "投融资服务", "基金＋银行＋创投协同，全年融资约9700余万元"),
    ("05", "区域与长三角协同", "校地合作、分园联动、产业资源共享"),
    ("06", "国际化合作", "北欧会客厅、出海沙龙与国际参访交流"),
]
for i, (no, t, b) in enumerate(feats):
    x = 0.5 + (i % 3) * 4.2
    y = 1.35 + (i // 3) * 2.8
    rect(s, x, y, 4.0, 2.5, CARD)
    text(s, x + 0.25, y + 0.3, 3.5, 0.45, [(no, 20, GOLD, True)])
    text(s, x + 0.25, y + 0.85, 3.5, 0.45, [(t, 16, NAVY, True)])
    text(s, x + 0.25, y + 1.45, 3.5, 0.8, [(b, 12, DARK, False)])

# ============================================================================
# 11 高校创业教育
# ============================================================================
s = add_slide(WHITE)
header(s, "特色工作 · 助力复旦大学创业教育", 11)
card(s, 0.55, 1.3, 6.1, 5.4, "主要抓手", [
    "双创课程：“复・创课堂”、企业沙龙、理论＋实践课程",
    "马兰花计划：GYB创业意识激发期培训班顺利举办",
    "创业见习与实训：强化学生项目筛选与创业实践",
    "品牌赛事：承办/协同复旦科创大赛、复旦之星等",
    "场地政策：持续向师生、校友创业组织提供租金减免",
    "2025年向49家创业组织提供租金减免等场地支持",
])
picture(s, "alumni1.jpeg", 6.9, 1.3, 5.8, 3.2)
picture(s, "alumni2.jpeg", 6.9, 4.7, 5.8, 2.0)

# ============================================================================
# 12 成果转化
# ============================================================================
s = add_slide(WHITE)
header(s, "特色工作 · 推动高校技术转移转化落地", 12)
text(s, 0.55, 1.2, 7.3, 2.2, [
    ("建设概念验证与联合孵化能力", 16, NAVY, True, 8),
    ("复煜概念验证中心2025年8月投入运营，从60余个项目中遴选4个进入验证阶段，提供资金、载体、场景、科研、投融资与运营协助。", 13, DARK, False, 8),
    ("累计对接复旦大学16个院系，新增入库成果21项，推动企业依托高校开展成果转化12项；深度服务11家复旦科创衍生企业。", 13, DARK, False),
])
card(s, 0.55, 3.7, 7.3, 3.0, "案例：波达医疗", [
    "复旦信息学院教师团队创办，超快超声成像产业化",
    "获“3310”B类、海聚英才大赛二等奖",
    "对接融资资源，获飞图创投数千万元融资，估值破亿",
    "获批上海市2025年度关键技术研发计划项目",
], accent=RED)
picture(s, "tech_transfer1.jpeg", 8.1, 1.25, 4.6, 2.7)
picture(s, "tech_transfer2.jpeg", 8.1, 4.2, 4.6, 2.5)

# ============================================================================
# 13 就业与融资
# ============================================================================
s = add_slide(WHITE)
header(s, "特色工作 · 就业促进与投融资服务", 13)
card(s, 0.5, 1.3, 6.1, 2.7, "就业促进", [
    "在孵组织在岗324人，新增就业118人",
    "开展校园招聘、企业实训与见习对接",
    "服务高校毕业生与青年创业者就业创业",
])
card(s, 6.85, 1.3, 5.9, 2.7, "投融资服务", [
    "投融资支持累计约9700余万元",
    "联动大创基金、科创母基金、前孵化资金",
    "合作创投与银行机构持续拓宽融资渠道",
], accent=GOLD)
picture(s, "recruit1.jpeg", 0.5, 4.3, 6.1, 2.5)
picture(s, "finance1.jpeg", 6.85, 4.3, 5.9, 2.5)

# ============================================================================
# 14 区域协同与国际化
# ============================================================================
s = add_slide(WHITE)
header(s, "特色工作 · 区域协同与国际化合作", 14)
card(s, 0.5, 1.3, 6.1, 3.0, "区域与开放协同", [
    "校地合作服务杨浦重点产业与成果落地",
    "拨投结合模式助力益臻新能源等项目",
    "基地企业区级税收贡献1.37亿元",
    "链接技术转移中心、张江磁谷等创新主体",
])
card(s, 6.85, 1.3, 5.9, 3.0, "国际化合作", [
    "参与发起“北欧创新国际会客厅”",
    "举办出海沙龙、学术交流等活动",
    "接待巴西、俄罗斯等国际参访团",
    "促进中欧创新资源与青年创业合作",
], accent=GOLD)
picture(s, "yangtze1.jpeg", 0.5, 4.55, 6.1, 2.2)
picture(s, "intl1.jpeg", 6.85, 4.55, 5.9, 2.2)

# ============================================================================
# 15 品牌赛事
# ============================================================================
s = add_slide(WHITE)
header(s, "特色品牌 · 复旦科创大赛", 15)
text(s, 0.55, 1.25, 12.2, 1.5, [
    ("由复旦大学、虹口区人民政府、锦江国际集团共同指导，基地作为主要承办单位，承担赛事组织、项目遴选、辅导培训、评审协调、资源对接及赛后孵化。", 14, DARK, False, 8),
    ("吸引全国456个项目报名，覆盖信息技术、集成电路、生命健康、新材料等领域；20支团队进入总决赛，多个项目达成投融资或转化合作意向。", 14, DARK, False),
])
stats = [("报名项目", "456个"), ("决赛团队", "20支"), ("重点赛道", "4大领域"), ("角色定位", "承办孵化")]
for i, (k, v) in enumerate(stats):
    x = 0.55 + i * 3.15
    rect(s, x, 3.2, 3.0, 1.7, CARD)
    text(s, x + 0.2, 3.4, 2.6, 0.4, [(k, 12, MUTED, False)])
    text(s, x + 0.2, 3.95, 2.6, 0.6, [(v, 24, NAVY, True)])
text(s, 0.55, 5.3, 12.2, 1.4, [
    ("以赛引才、以赛选才、以赛聚才，推动创新链与产业链深度对接，为国家高质量发展注入高校智慧与青年动能。", 14, DARK, False),
    ("同步协同“复旦之星”“创·在上海”“杨浦科创之星”等赛事，完善高层次人才与优质项目集聚平台。", 14, DARK, False),
])

# ============================================================================
# 16 典型案例
# ============================================================================
s = add_slide(WHITE)
header(s, "孵化典型案例", 16)
rect(s, 0.5, 1.25, 7.5, 5.5, CARD)
text(s, 0.8, 1.5, 7.0, 0.5, [("波达医疗 · 高校成果转化标杆", 20, NAVY, True)])
text(s, 0.8, 2.2, 7.0, 4.2, [
    ("企业定位：复旦大学信息学院教师团队创办，聚焦超快超声成像设备与技术产业化。", 13, DARK, False, 8),
    ("孵化服务：创业辅导、3310申报、载体落地、租金与社保政策、投融资对接。", 13, DARK, False, 8),
    ("成长成效：海聚英才大赛二等奖；获飞图创投数千万元融资，估值破亿；获批上海市2025年度关键技术研发计划项目。", 13, DARK, False, 8),
    ("同期典型：眸深智能获全国颠覆性技术创新大赛最高奖；本科生团队获中国国际大学生创新大赛金奖；音书科技获天使基金优秀雏鹰企业。", 13, DARK, False, 8),
])
picture(s, "tech_transfer1.jpeg", 8.3, 1.25, 4.4, 2.6)
picture(s, "tech_transfer2.jpeg", 8.3, 4.05, 4.4, 2.7)

# ============================================================================
# 17 综合效益
# ============================================================================
s = add_slide(WHITE)
header(s, "综合效益", 17)
eff = [
    ("经济效益", ["新增注册企业190家（科技型143家）", "高企（含复审）21家 / 小巨人2家 / 3310企业3家", "知识产权318项（发明专利42项）", "区级税收贡献1.37亿元"]),
    ("社会效益", ["成果转化12项，新增导师16人", "培训11场、活动37场，赛事推荐100余项", "市载体绩效：科技园优良、众创空间优秀", "杨浦区科技园区考评优秀"]),
    ("绿色发展", ["景观照明与公共空间品质提升", "综合更新项目获批张江专项", "获评区级“无废楼宇”", "持续优化载体运营环境"]),
]
for i, (t, lines) in enumerate(eff):
    card(s, 0.5 + i * 4.2, 1.4, 4.0, 5.2, t, lines, accent=[NAVY, BLUE, GOLD][i])

# ============================================================================
# 18 封底
# ============================================================================
s = add_slide(WHITE)
picture(s, "closing_photo.jpeg", 6.4, 0, 7.0, 7.5)
rect(s, 0, 0, 7.0, 7.5, NAVY)
rect(s, 0, 6.85, SW_IN, 0.65, NAVY_DK)
text(s, 0.8, 2.3, 5.8, 2.4, [
    ("育高校创业新苗，", 30, WHITE, True, 8),
    ("助技术成果转化！", 30, WHITE, True, 16),
    ("复旦大学国家大学科技园", 16, GOLD, False, 6),
    ("创业孵化基地 · 2026年工作报告", 14, RGBColor(0xA9, 0xC0, 0xD6), False),
])
text(s, 0.8, 7.0, 5.8, 0.35, [("感谢聆听  ·  2026.06", 13, WHITE, False)], anchor=MSO_ANCHOR.MIDDLE)

prs.save(OUT)
print(f"已生成：{OUT}")
