#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
创智汇年度活动运营方案（双会并入融合版 · 甲方汇报版）PPT 生成脚本

设计模板（Kimi K3 汇报模板）：
  · 米白底 + 深海军蓝 + 香槟金点缀
  · 每页配统一风格插图（assets/chuangzhihui/）
  · 闭环叙事：背景 → 主线 → 框架 → 日历 → 节奏 → 主题 → 单场打法 →
    转化闭环 → 交付 → 收费 → 案例 → 结论（呼应开篇）

运行：python3 scripts/build_chuangzhihui_annual_deck.py
输出：output/创智汇年度活动运营方案-双会并入融合版-汇报版.pptx
"""
import copy
import os

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------- 设计规范
NAVY = RGBColor(0x0E, 0x22, 0x40)       # 深海军蓝（主色）
NAVY_DEEP = RGBColor(0x0A, 0x19, 0x30)  # 封面深底
NAVY_SOFT = RGBColor(0x1B, 0x3A, 0x6B)  # 浅一档蓝
GOLD = RGBColor(0xC9, 0xA2, 0x27)       # 香槟金（点缀）
GOLD_DEEP = RGBColor(0xA8, 0x86, 0x1B)
GOLD_PALE = RGBColor(0xF4, 0xEB, 0xD3)  # 金色浅底
CREAM = RGBColor(0xFA, 0xF6, 0xEE)      # 米白底
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x24, 0x30, 0x44)        # 正文深色
GRAY = RGBColor(0x6B, 0x72, 0x80)       # 次要灰
LINE = RGBColor(0xE7, 0xDD, 0xC6)       # 卡片描边
BLUE_TAG = RGBColor(0x2E, 0x5A, 0xA8)   # WAIC 线蓝

FONT = "微软雅黑"

PAGE_W = Inches(13.333)
PAGE_H = Inches(7.5)

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ASSETS = os.path.join(ROOT, "assets", "chuangzhihui")
CROP_DIR = os.path.join(ASSETS, "_crop")
OUT = os.path.join(ROOT, "output", "创智汇年度活动运营方案-双会并入融合版-汇报版.pptx")

FOOTER = "上海创智汇 × 同浦汇 · 创智汇年度活动运营方案 · 2026.08—2027.07"
TOTAL = 18

os.makedirs(CROP_DIR, exist_ok=True)
os.makedirs(os.path.dirname(OUT), exist_ok=True)


# ---------------------------------------------------------------- 基础工具
def crop_to(path, ratio, tag):
    """把图片中心裁剪到目标宽高比，返回裁剪后路径。"""
    dst = os.path.join(CROP_DIR, f"{tag}.png")
    if os.path.exists(dst):
        return dst
    im = Image.open(path).convert("RGB")
    w, h = im.size
    cur = w / h
    if cur > ratio:
        nw = int(h * ratio)
        x0 = (w - nw) // 2
        im = im.crop((x0, 0, x0 + nw, h))
    else:
        nh = int(w / ratio)
        y0 = (h - nh) // 2
        im = im.crop((0, y0, w, y0 + nh))
    im.save(dst)
    return dst


def img(name):
    return os.path.join(ASSETS, name)


def _style_run(run, size, color, bold=False, font=FONT, spacing=None, italic=False):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = font
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", font)
    if spacing is not None:
        rPr.set("spc", str(int(spacing * 100)))


def add_text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.12, wrap=True):
    """lines: [ [(text,size,color,bold), ...run], ...paragraph] 或简化 [(text,size,color,bold)]"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for para in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = line_spacing
        runs = para if isinstance(para, list) else [para]
        for item in runs:
            text, size, color, bold = item[0], item[1], item[2], item[3]
            r = p.add_run()
            r.text = text
            _style_run(r, size, color, bold)
    return tb


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, radius=None,
             shadow=False, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shadow:
        _add_shadow(sp)
    sp.text_frame.word_wrap = True
    return sp


def _add_shadow(shape, blur=110000, dist=38000, alpha=16000):
    spPr = shape._element.spPr
    old = spPr.find(qn("a:effectLst"))
    if old is not None:
        spPr.remove(old)
    eff = spPr.makeelement(qn("a:effectLst"), {})
    shdw = spPr.makeelement(qn("a:outerShdw"), {
        "blurRad": str(blur), "dist": str(dist), "dir": "5400000", "rotWithShape": "0"})
    clr = spPr.makeelement(qn("a:srgbClr"), {"val": "0E2240"})
    a = spPr.makeelement(qn("a:alpha"), {"val": str(alpha)})
    clr.append(a)
    shdw.append(clr)
    eff.append(shdw)
    spPr.append(eff)


def set_alpha(shape, opacity_pct):
    """给纯色填充加透明度。opacity_pct: 0-100（100 为不透明）。"""
    sf = shape.fill._xPr.find(qn("a:solidFill"))
    if sf is None:
        return
    clr = sf.find(qn("a:srgbClr"))
    if clr is None:
        return
    a = clr.makeelement(qn("a:alpha"), {"val": str(int(opacity_pct * 1000))})
    clr.append(a)


def add_pic(slide, path, x, y, w, h, border=GOLD, border_w=1.0, shadow=True):
    pic = slide.shapes.add_picture(path, x, y, w, h)
    if border is not None:
        pic.line.color.rgb = border
        pic.line.width = Pt(border_w)
    pic.shadow.inherit = False
    if shadow:
        _add_shadow(pic)
    return pic


def shape_text(sp, lines, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1):
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    first = True
    for para in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = line_spacing
        runs = para if isinstance(para, list) else [para]
        for item in runs:
            text, size, color, bold = item[0], item[1], item[2], item[3]
            r = p.add_run()
            r.text = text
            _style_run(r, size, color, bold)


def add_bg(slide, color=CREAM):
    add_rect(slide, 0, 0, PAGE_W, PAGE_H, fill=color, shape=MSO_SHAPE.RECTANGLE)


def add_footer(slide, n, dark=False):
    c = GOLD_PALE if dark else GRAY
    ln = GOLD if dark else LINE
    add_rect(slide, Inches(0.42), Inches(7.12), Inches(12.49), Pt(0.9),
             fill=ln, shape=MSO_SHAPE.RECTANGLE)
    add_text(slide, Inches(0.42), Inches(7.18), Inches(9.6), Inches(0.26),
             [(FOOTER, 8.5, c, False)])
    add_text(slide, Inches(11.3), Inches(7.18), Inches(1.61), Inches(0.26),
             [(f"{n:02d} / {TOTAL}", 8.5, c, False)], align=PP_ALIGN.RIGHT)


def add_header(slide, tag, title, subtitle=None):
    add_rect(slide, Inches(0.42), Inches(0.42), Inches(0.14), Inches(0.14),
             fill=GOLD, shape=MSO_SHAPE.RECTANGLE)
    add_text(slide, Inches(0.66), Inches(0.33), Inches(8.0), Inches(0.3),
             [(tag, 12, GOLD_DEEP, True)])
    add_text(slide, Inches(0.42), Inches(0.62), Inches(10.6), Inches(0.55),
             [(title, 24, NAVY, True)])
    if subtitle:
        add_text(slide, Inches(0.42), Inches(1.18), Inches(11.9), Inches(0.3),
                 [(subtitle, 11, GRAY, False)])
    add_rect(slide, Inches(0.42), Inches(1.56), Inches(1.6), Pt(2.2),
             fill=GOLD, shape=MSO_SHAPE.RECTANGLE)


def chip(slide, x, y, w, h, text, fill=WHITE, color=NAVY, size=9.5, bold=True,
         line=LINE, radius=0.5, shadow=False):
    sp = add_rect(slide, x, y, w, h, fill=fill, line=line, radius=radius, shadow=shadow)
    shape_text(sp, [(text, size, color, bold)])
    return sp


def badge(slide, x, y, d, text, fill=GOLD, color=WHITE, size=13, bold=True):
    sp = add_rect(slide, x, y, d, d, fill=fill, shape=MSO_SHAPE.OVAL)
    shape_text(sp, [(text, size, color, bold)])
    return sp


def add_chevron(slide, x, y, w, h, text_lines, fill=NAVY, first=False, flip_h=False):
    shape_type = MSO_SHAPE.PENTAGON if first else MSO_SHAPE.CHEVRON
    sp = add_rect(slide, x, y, w, h, fill=fill, shape=shape_type)
    if flip_h:
        sp._element.spPr.find(qn("a:xfrm")).set("flipH", "1")
    shape_text(sp, text_lines)
    return sp


def style_table(gf, col_widths, header_fill=NAVY):
    tbl = gf._element.graphic.graphicData.tbl
    tblPr = tbl.find(qn("a:tblPr"))
    if tblPr is not None:
        tblPr.set("firstRow", "0")
        tblPr.set("bandRow", "0")
    table = gf.table
    total = sum(col_widths)
    for i, cwd in enumerate(col_widths):
        table.columns[i].width = Emu(int(cwd / total * gf.width))
    return table


def set_cell(cell, text, size=10, color=INK, bold=False, fill=None,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE
    cell.vertical_anchor = anchor
    cell.margin_left = cell.margin_right = Inches(0.045)
    cell.margin_top = cell.margin_bottom = Inches(0.012)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = 1.05
    r = p.add_run()
    r.text = text
    _style_run(r, size, color, bold)


# ---------------------------------------------------------------- 数据
CAL_H1 = [
    ("P1", "创智汇项目推广日①·中介与渠道专场", "8月上旬", "项目推广日", "≤30人", "L0", "中介渠道带看与项目推介"),
    ("P2", "创智汇项目推广日②·五大行与金融", "8月中旬", "项目推广日", "≤30人", "L0", "金融资源链接与客户导入"),
    ("P3", "创智汇项目推广日③·投研机构专场", "8月中旬", "项目推广日", "≤30人", "L0", "投研视角放大项目影响力"),
    ("P4", "创智汇项目推广日④·政府部门与载体", "8月下旬", "项目推广日", "≤30人", "L0", "政策与载体协同落地"),
    ("E1", "WAIC成果承接·创智汇开放交流日", "8月下旬", "开放交流日", "≤30人", "L1", "衔接WAIC议题与园区看场"),
    ("A2", "Agent智能体搭建一日营", "9月上旬", "实训营", "≤30人", "L2", "Agent初创落位"),
    ("C1", "ChinaJoy主题·数字娱乐与IP授权对接", "9月中旬", "对接会", "≤30人", "L2", "CJ赛道企业/IP方入园"),
    ("F1", "出海专题活动（另议·另计价）", "档期另议", "专题（另议）", "≤30人", "另计价", "出海企业服务·不占年度包"),
    ("A3", "多模态智能体工作坊", "10月上旬", "工作坊", "≤30人", "L2", "多模态团队入驻"),
    ("E4", "ChinaJoy主题日·数字娱乐生态", "10月中旬", "主题日", "≤30人", "L2", "CJ生态企业集中看场"),
    ("B2", "AI治理与可信智能体沙龙", "10月下旬", "沙龙", "≤30人", "L2", "合规型企业信任导入"),
    ("A4", "火山引擎×算力Infra实务营", "11月上旬", "厂商联训", "≤30人", "L2", "高算力企业定向看场"),
    ("B4", "创新券·算力券·模型券实务沙龙", "11月中旬", "实务沙龙", "≤30人", "L2", "用券企业聚集"),
    ("L1", "领事到访接待（另计价）", "档期另议", "外事接待", "另议", "另计价", "与挂牌分属不同性质·分开计价"),
    ("A5", "具身智能空间交互体验日", "12月上旬", "体验日", "≤30人", "L3", "具身/机器人团队看场"),
]

CAL_H2 = [
    ("F4", "创智汇AI年度Demo Day", "12月中旬", "路演日", "≤30人", "L4", "集中签约洽谈与媒体背书"),
    ("D1", "AIGC微短剧制片特训", "1月上旬", "特训营", "≤30人", "L2", "厂牌/工作室入驻"),
    ("A6", "AI营销Agent实战营", "1月中旬", "实战营", "≤30人", "L2", "营销科技公司看场"),
    ("B1", "YOUNG立方×智能伙伴政策沙龙", "2月中旬", "政策沙龙", "≤30人", "L2", "内容/AI企业导入"),
    ("B3", "高企认定冲刺（AI企业专场）", "2月下旬", "辅导会", "≤30人", "L2", "待认定企业带政策入驻"),
    ("A7", "OPC超级个体黑客松（春）", "3月上旬", "黑客松", "≤30人", "L2", "获奖团队优先谈单元"),
    ("C2", "高校成果转化·AI for Science对接日", "3月下旬", "对接日", "≤30人", "L2", "成果公司/实验室落户"),
    ("F3", "通往AGI季度圆桌", "3月下旬", "闭门圆桌", "≤30人", "L2", "研究型/模型团队"),
    ("D2", "ChinaJoy主题·创作者内容首发", "4月中旬", "发布会", "≤30人", "L2", "数字娱乐/内容品牌问询"),
    ("C3", "ChinaJoy主题·游戏周边与潮玩供应链", "4月下旬", "对接会", "≤30人", "L2", "CJ供应链/周边企业展位"),
    ("C4", "专精特新·AI应用培育路演", "5月中旬", "路演", "≤30人", "L2", "成长型AI企业补位"),
    ("D3", "ChinaJoy主题·数字娱乐市集体验日", "5月下旬", "体验日", "≤30人", "L2", "摊主/工作室升级固定展位"),
    ("C5", "ChinaJoy主题·IP联名与衍生品撮合", "6月中旬", "撮合会", "≤30人", "L2", "IP联名/衍生品团队入驻"),
    ("D4", "ChinaJoy主题·沉浸式数字联展", "6月下旬", "联展(短展期)", "≤30人/场", "L2", "闭幕洽谈集中转化"),
    ("D5", "创作者首发②·衔接WAIC 2027预热", "7月中旬", "发布会", "≤30人", "L4", "旺季补位·双大会预热"),
]

THEMES = {
    "A": dict(name="智能体与算力训练", count="6场", span="9月—次年3月",
              role="3F基本盘：营 / 坊 / 联训 / 体验 / 黑客松",
              res="对应 WAIC：Agent·多模态·算力·具身；ACT缩尺+线索抽样",
              events=[("A2", "Agent智能体搭建一日营", "9月上旬"),
                      ("A3", "多模态智能体工作坊", "10月上旬"),
                      ("A4", "火山引擎×算力Infra实务营", "11月上旬"),
                      ("A5", "具身智能空间交互体验日", "12月上旬"),
                      ("A6", "AI营销Agent实战营", "次年1月中旬"),
                      ("A7", "OPC超级个体黑客松（春）", "次年3月上旬")],
              point="Agent/多模态/算力/具身团队看场落位"),
    "B": dict(name="政策与治理沙龙", count="4场", span="10月—次年2月",
              role="政策 / 券务 / 资质，做成进园理由",
              res="对应 WAIC 治理与券务议题；服中心协同兑现",
              events=[("B2", "AI治理与可信智能体沙龙", "10月下旬"),
                      ("B4", "创新券·算力券·模型券实务沙龙", "11月中旬"),
                      ("B1", "YOUNG立方×智能伙伴政策沙龙", "次年2月中旬"),
                      ("B3", "高企认定冲刺（AI企业专场）", "次年2月下旬")],
              point="合规型企业信任导入、带政策入驻"),
    "C": dict(name="产业与ChinaJoy对接", count="5场", span="全年各月",
              role="高校成果 + CJ工作室 / 供应链 / IP",
              res="原潮玩产业集群场次合并为 ChinaJoy 主题内容",
              events=[("C1", "数字娱乐与IP授权对接会", "9月中旬"),
                      ("C2", "高校成果转化·AI for Science", "次年3月下旬"),
                      ("C3", "游戏周边与潮玩供应链对接", "次年4月下旬"),
                      ("C4", "专精特新·AI应用培育路演", "次年5月中旬"),
                      ("C5", "IP联名与衍生品撮合会", "次年6月中旬")],
              point="CJ赛道企业/IP方/成果公司入园"),
    "D": dict(name="AI内容与ChinaJoy人气", count="5场", span="1月—7月",
              role="5F人气引擎：特训 / 首发 / 市集 / 联展",
              res="对应 CJ 创作者经济 + WAIC AIGC；与官方活动错峰",
              events=[("D1", "AIGC微短剧制片特训", "次年1月上旬"),
                      ("D2", "创作者内容首发发布会", "次年4月中旬"),
                      ("D3", "数字娱乐市集体验日", "次年5月下旬"),
                      ("D4", "沉浸式数字联展（短展期）", "次年6月下旬"),
                      ("D5", "创作者首发②·WAIC预热", "次年7月中旬")],
              point="内容厂牌/工作室入驻、旺季补位预热"),
    "E": dict(name="启动月·项目推广与双会承接", count="6场", span="8月、10月",
              role="8月点火 + WAIC承接 + CJ主题日",
              res="原大场改为≤30人精准场；8月=项目推广日",
              events=[("P1", "推广日①·中介与渠道", "8月上旬"),
                      ("P2", "推广日②·五大行与金融", "8月中旬"),
                      ("P3", "推广日③·投研机构专场", "8月中旬"),
                      ("P4", "推广日④·政府部门与载体", "8月下旬"),
                      ("E1", "WAIC成果承接开放交流日", "8月下旬"),
                      ("E4", "ChinaJoy主题日·数字娱乐", "10月中旬")],
              point="渠道/金融/投研/政府四类资源点火"),
    "F": dict(name="收官·圆桌·另计价事项", count="2场+另计价", span="12月、次年3月",
              role="Demo Day 收口；出海 / 领事另案",
              res="出海另议另计价；领事到访与挂牌分属不同性质",
              events=[("F4", "创智汇AI年度Demo Day", "12月中旬"),
                      ("F3", "通往AGI季度圆桌", "次年3月下旬"),
                      ("F1", "出海专题活动（另议·另计价）", "档期另议"),
                      ("L1", "领事到访接待（另计价）", "档期另议")],
              point="集中签约洽谈与媒体背书、研究型客户"),
}


# ---------------------------------------------------------------- 页面构建
def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, PAGE_W, PAGE_H, fill=NAVY_DEEP, shape=MSO_SHAPE.RECTANGLE)
    # 全幅主视觉（16:9 与页面同比例）
    s.shapes.add_picture(img("cover_hero.png"), 0, 0, PAGE_W, PAGE_H)
    # 左侧压暗渐变（两段半透明叠加）
    ov1 = add_rect(s, 0, 0, Inches(8.6), PAGE_H, fill=NAVY_DEEP, shape=MSO_SHAPE.RECTANGLE)
    set_alpha(ov1, 62)
    ov2 = add_rect(s, 0, 0, Inches(6.4), PAGE_H, fill=NAVY_DEEP, shape=MSO_SHAPE.RECTANGLE)
    set_alpha(ov2, 55)
    # 顶部金色 kicker
    add_rect(s, Inches(0.55), Inches(0.62), Inches(0.5), Pt(2.6), fill=GOLD, shape=MSO_SHAPE.RECTANGLE)
    add_text(s, Inches(0.55), Inches(0.74), Inches(9.5), Inches(0.34),
             [("上海创智汇 × 同浦汇 ｜ 年度活动专项 · 甲方汇报版", 12.5, GOLD, True)])
    # 主标题
    add_text(s, Inches(0.55), Inches(1.28), Inches(9.8), Inches(1.9),
             [("创智汇年度活动运营方案", 41, WHITE, True),
              ("双会并入 · 30场活动怎么做", 25, GOLD, True)],
             line_spacing=1.22)
    # 副题
    add_text(s, Inches(0.55), Inches(3.28), Inches(9.5), Inches(0.4),
             [("2026.08 — 2027.07　｜　同步衔接 WAIC 与 ChinaJoy　｜　把双会热度，变成园区的一年", 12.5,
               GOLD_PALE, False)])
    # 关键数字条
    stats = [("30场", "全年活动总量"), ("≤30人", "每场到场人数"), ("约50%", "负责人占比承兑"), ("600+", "全年触达人次")]
    cx = 0.55
    for num, lab in stats:
        card = add_rect(s, Inches(cx), Inches(4.05), Inches(2.16), Inches(0.92),
                        fill=NAVY, radius=0.16, shadow=True)
        set_alpha(card, 78)
        shape_text(card, [[(num, 19, GOLD, True)], [(lab, 10, GOLD_PALE, False)]], line_spacing=1.05)
        cx += 2.36
    # 底部信息
    add_text(s, Inches(0.55), Inches(5.42), Inches(11.0), Inches(1.1),
             [("汇报对象：上海创智汇（园区）　·　汇报与执行：同浦汇", 11, WHITE, False),
              ("主体运营：上海市云计算创新基地（国家级孵化器）　·　学术支持：复旦大学住房政策研究中心", 10, GOLD_PALE, False),
              ("载体支持：杨浦区科技企业联合会 · 科技企业服务中心", 10, GOLD_PALE, False)],
             line_spacing=1.45)
    add_text(s, Inches(0.55), Inches(6.95), Inches(8.0), Inches(0.3),
             [("2026年8月 · 上海", 10, GOLD_PALE, False)])


def slide_toc(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_text(s, Inches(0.42), Inches(0.5), Inches(6.0), Inches(0.4),
             [("CONTENTS", 12, GOLD_DEEP, True)])
    add_text(s, Inches(0.42), Inches(0.78), Inches(6.0), Inches(0.62),
             [("汇报导览", 30, NAVY, True)])
    add_rect(s, Inches(0.42), Inches(1.5), Inches(1.6), Pt(2.2), fill=GOLD, shape=MSO_SHAPE.RECTANGLE)

    entries = [
        ("01", "背景与总体思路", "为什么做：双会热度需要承接地，一条主线、六大主题", "P.03"),
        ("02", "年度日历与运营节奏", "30场排期总表、四阶段节奏、双会资源导入摘要", "P.06"),
        ("03", "六大主题活动详述", "智能体算力 · 政策治理 · 产业CJ · 内容人气 · 启动承接 · 收官", "P.09"),
        ("04", "单场执行标准", "每场从 T-14 筹备到 T+7 回访的同一套打法", "P.12"),
        ("05", "转化闭环 · 交付 · 收费", "到场→建群→看场→促成→回流；交付清单与付款节点", "P.13"),
        ("06", "往期案例与一页纸结论", "已跑通的两类活动案例；我们交付 / 园区协同 / 商务口径", "P.16"),
    ]
    y = 1.8
    for num, title, desc, page in entries:
        card = add_rect(s, Inches(0.42), Inches(y), Inches(8.55), Inches(0.7),
                        fill=WHITE, line=LINE, radius=0.14, shadow=True)
        badge(s, Inches(0.62), Inches(y + 0.15), Inches(0.4), num, fill=NAVY, size=11)
        add_text(s, Inches(1.22), Inches(y + 0.07), Inches(6.2), Inches(0.3),
                 [(title, 13.5, NAVY, True)])
        add_text(s, Inches(1.22), Inches(y + 0.37), Inches(6.6), Inches(0.26),
                 [(desc, 9.5, GRAY, False)])
        add_text(s, Inches(7.9), Inches(y + 0.19), Inches(0.95), Inches(0.3),
                 [(page, 11, GOLD_DEEP, True)], align=PP_ALIGN.RIGHT)
        y += 0.82
    # 阅读主线（闭环提示）
    add_text(s, Inches(0.42), Inches(6.78), Inches(8.7), Inches(0.3),
             [[("阅读主线：", 10, GOLD_DEEP, True),
               ("为什么做 → 怎么排 → 怎么落地 → 怎么转化收费 → 凭什么做好，首尾呼应成闭环", 10, GRAY, False)]])
    # 右侧插图
    pic = crop_to(img("toc_roadmap.png"), 3.45 / 5.6, "toc_rail")
    add_pic(s, pic, Inches(9.25), Inches(0.42), Inches(3.45), Inches(5.6))
    chip(s, Inches(9.25), Inches(6.2), Inches(3.45), Inches(0.44),
         "双会是资源库 · 30场是定义核", fill=NAVY, color=GOLD_PALE, size=10.5, line=None, radius=0.3)
    add_footer(s, n)


def slide_background(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "01 · 背景与判断", "大会闭幕之后，热度需要一个承接地",
               "主题响应：上海创智汇 · AI+数字内容无界共创港")
    cards = [
        ("双会已闭幕，热度在账上", GOLD,
         [("WAIC 2026 已于 7 月闭幕，ChinaJoy 以「与AI同游」同步释放数字内容与具身线索；", 10.5),
          ("大会积累的嘉宾、议题与可触达线索若只停留在会期，对园区的带动就到此为止。", 10.5)]),
        ("园区在爬坡，缺一整年经营", NAVY_SOFT,
         [("3F 智能体与算力组团仍有单元虚位，5F 内容与 IP 组团需要持续人气；", 10.5),
          ("两个楼层缺的不是一次性活动，而是一整年有节奏的经营；租金高于周边均价，更需经营佐证价值。", 10.5)]),
        ("我们的判断", GOLD_DEEP,
         [("双会是资源库，30场是定义核——2026.08 启动、2027.07 收官，每月两到三场、精而准，", 10.5),
          ("落点放在看场、建群与销售促成（怎么做 → 见 P.04 总体思路）。", 10.5)]),
    ]
    y = 1.86
    for title, accent, lines in cards:
        card = add_rect(s, Inches(0.42), Inches(y), Inches(8.85), Inches(1.28),
                        fill=WHITE, line=LINE, radius=0.1, shadow=True)
        add_rect(s, Inches(0.42), Inches(y), Inches(0.09), Inches(1.28), fill=accent, shape=MSO_SHAPE.RECTANGLE)
        add_text(s, Inches(0.72), Inches(y + 0.12), Inches(8.3), Inches(0.3),
                 [(title, 13, NAVY, True)])
        paras = [[(t, size, INK, False)] for t, size in lines]
        add_text(s, Inches(0.72), Inches(y + 0.46), Inches(8.35), Inches(0.75), paras, line_spacing=1.28)
        y += 1.44
    # 关键数字条
    stats = [("30场", "覆盖双会主要招商赛道"), ("12个月", "每月2–3场不断档"),
             ("≤30人/场", "全年约600人次+可核验"), ("约50%", "负责人占比适度承兑")]
    cx = 0.42
    for num, lab in stats:
        card = add_rect(s, Inches(cx), Inches(6.28), Inches(2.13), Inches(0.66),
                        fill=NAVY, radius=0.14, shadow=True)
        shape_text(card, [[(num + "  ", 14, GOLD, True), (lab, 8.5, GOLD_PALE, False)]])
        cx += 2.24
    # 右侧插图
    pic = crop_to(img("bg_funnel.png"), 3.42 / 5.28, "bg_rail")
    add_pic(s, pic, Inches(9.49), Inches(0.42), Inches(3.42), Inches(5.28))
    chip(s, Inches(9.49), Inches(5.86), Inches(3.42), Inches(0.44),
         "热度不止于会期", fill=GOLD_PALE, color=GOLD_DEEP, size=10.5, line=None, radius=0.3)
    add_footer(s, n)


def slide_mainline(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "01 · 总体思路", "一条主线：把双会的热度，变成园区的一年",
               "双会资源库 → 六大主题×30场 → 看场·建群·销售促成（转化闭环详见 P.13）")
    # 三段主流程
    steps = [
        ("双会资源库", "WAIC → 3F 能力线\nChinaJoy → 5F 内容线\n具身 / AIGC 双会叠加带", NAVY),
        ("六大主题 × 30场", "每月 2–3 场不断档\n每场≤30人 · 负责人约50%\n只取招商贴近的资源切片", NAVY_SOFT),
        ("园区的一年", "固定看场 + 当场建群\n同浦汇带客 · 园区销售\n签到/建群/回访可归档", GOLD_DEEP),
    ]
    x = 0.42
    w = 3.94
    for i, (t, body, fill) in enumerate(steps):
        sp = add_chevron(s, Inches(x), Inches(1.86), Inches(w), Inches(1.5),
                         [[(t, 14.5, WHITE, True)]] + [[(ln, 9.5, GOLD_PALE, False)] for ln in body.split("\n")],
                         fill=fill, first=(i == 0))
        x += w + 0.14
    # 主视觉横幅
    pic = crop_to(img("mainline_merge.png"), 12.49 / 2.5, "mainline_banner")
    add_pic(s, pic, Inches(0.42), Inches(3.56), Inches(12.49), Inches(2.5))
    # 双收益卡片
    c1 = add_rect(s, Inches(0.42), Inches(6.22), Inches(6.14), Inches(0.78),
                  fill=WHITE, line=LINE, radius=0.12, shadow=True)
    add_text(s, Inches(0.62), Inches(6.3), Inches(5.9), Inches(0.62),
             [[("园区得到什么：", 10.5, GOLD_DEEP, True),
               ("议题变日常 · 空间有主题 · 考核有支撑 · 国家级孵化器背书可感知", 10, INK, False)]],
             line_spacing=1.25)
    c2 = add_rect(s, Inches(6.77), Inches(6.22), Inches(6.14), Inches(0.78),
                  fill=WHITE, line=LINE, radius=0.12, shadow=True)
    add_text(s, Inches(6.97), Inches(6.3), Inches(5.9), Inches(0.62),
             [[("招商得到什么：", 10.5, GOLD_DEEP, True),
               ("客群更准 · 动作标准化 · 政策券务打包降门槛 · 不做租赁对赌", 10, INK, False)]],
             line_spacing=1.25)
    add_footer(s, n)


def slide_framework(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "01 · 总体框架", "六大主题，各管一段，互相导流",
               "并入四层资源叠加：L0 点火 → L1 承接 → L2 主线 → L3 / L4 收口（逐场详述见 P.09—P.11）")
    order = ["A", "B", "C", "D", "E", "F"]
    pos = [(0.42, 1.86), (5.05, 1.86), (0.42, 3.36), (5.05, 3.36), (0.42, 4.86), (5.05, 4.86)]
    w, h = 4.45, 1.36
    for key, (x, y) in zip(order, pos):
        t = THEMES[key]
        add_rect(s, Inches(x), Inches(y), Inches(w), Inches(h), fill=WHITE, line=LINE, radius=0.1, shadow=True)
        badge(s, Inches(x + 0.16), Inches(y + 0.16), Inches(0.44), key, fill=GOLD if key in "EF" else NAVY, size=13)
        add_text(s, Inches(x + 0.74), Inches(y + 0.1), Inches(w - 0.9), Inches(0.3),
                 [[(t["name"], 12, NAVY, True), ("　" + t["count"], 10.5, GOLD_DEEP, True)]])
        add_text(s, Inches(x + 0.74), Inches(y + 0.42), Inches(w - 0.9), Inches(0.26),
                 [(t["role"], 9, GRAY, False)])
        add_text(s, Inches(x + 0.16), Inches(y + 0.74), Inches(w - 0.32), Inches(0.26),
                 [[("时间分布　", 8.5, GOLD_DEEP, True), (t["span"], 9.5, INK, False)]])
        add_text(s, Inches(x + 0.16), Inches(y + 1.0), Inches(w - 0.32), Inches(0.26),
                 [[("落点　　　", 8.5, GOLD_DEEP, True), (t["point"], 9.5, INK, False)]])
    # 右侧插图
    pic = crop_to(img("framework_hex.png"), 3.42 / 4.9, "hex_rail")
    add_pic(s, pic, Inches(9.49), Inches(0.42), Inches(3.42), Inches(4.9))
    chip(s, Inches(9.49), Inches(5.48), Inches(3.42), Inches(0.74),
         "28场常规活动 + 2项另计价\n覆盖双会主要招商赛道", fill=NAVY, color=GOLD_PALE, size=10, line=None, radius=0.2)
    add_footer(s, n)


def _calendar_slide(prs, n, rows, half, span_desc, rail_img, rail_chips, note):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "02 · 年度日历", f"30场排期总表（{half}）",
               f"{span_desc}　｜　每场≤30人 · 负责人约50%")
    # 表格
    col_w = [0.62, 2.75, 1.02, 1.08, 0.78, 0.72, 2.17]
    gf = s.shapes.add_table(len(rows) + 1, 7, Inches(0.42), Inches(1.78),
                            Inches(9.14), Inches(5.2))
    table = style_table(gf, col_w)
    table.rows[0].height = Inches(0.34)
    headers = ["编号", "活动", "档期", "形式", "人数", "层次", "招商落点"]
    for c, htext in enumerate(headers):
        set_cell(table.cell(0, c), htext, size=10, color=GOLD_PALE, bold=True, fill=NAVY)
    for i, row in enumerate(rows):
        table.rows[i + 1].height = Inches(0.319)
        fill = WHITE if i % 2 == 0 else GOLD_PALE
        for c, val in enumerate(row):
            bold = c == 0
            color = NAVY if c == 0 else INK
            align = PP_ALIGN.LEFT if c in (1, 6) else PP_ALIGN.CENTER
            set_cell(table.cell(i + 1, c), val, size=9, color=color, bold=bold,
                     fill=fill, align=align)
    # 右栏插图 + 要点
    pic = crop_to(img(rail_img), 3.42 / 2.3, f"cal_{half}")
    add_pic(s, pic, Inches(9.49), Inches(0.42), Inches(3.42), Inches(2.3))
    y = 2.9
    for c in rail_chips:
        chip(s, Inches(9.49), Inches(y), Inches(3.42), Inches(0.42), c,
             fill=WHITE, color=NAVY, size=9.5, line=LINE, radius=0.3, shadow=True)
        y += 0.54
    note_card = add_rect(s, Inches(9.49), Inches(y + 0.04), Inches(3.42), Inches(6.94 - y - 0.04),
                         fill=GOLD_PALE, radius=0.1)
    add_text(s, Inches(9.65), Inches(y + 0.14), Inches(3.12), Inches(6.8 - y),
             [[("注　", 9.5, GOLD_DEEP, True), (note, 9, INK, False)]], line_spacing=1.3)
    add_footer(s, n)


def slide_rhythm(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "02 · 运营节奏与双会导入", "四个阶段各有重心，双会各管一楼",
               "L0点火 → L1承接 → L2主线 → L3/L4收口；只取招商贴近可调用规模与优先策略")
    stages = [
        ("启动期｜8—10月", ["四场项目推广日点火", "E1 承接 WAIC 余热", "Agent营 + CJ IP对接", "CJ主题日 + 治理沙龙"], NAVY),
        ("攻坚期｜11—12月", ["算力联训 + 三券沙龙", "具身双会叠加体验", "Demo Day 精选签约", "领事另计价另案（如需）"], NAVY_SOFT),
        ("深耕期｜1—3月", ["AIGC特训 / 营销Agent", "政策沙龙 + 高企冲刺", "黑客松 + 高校成果日", "AGI圆桌研究型客户"], BLUE_TAG),
        ("收获与预热｜4—7月", ["CJ创作者首发 / 供应链", "市集转正 / IP联名 / 联展", "专精特新AI路演补位", "D5 预热下届双会"], GOLD_DEEP),
    ]
    pos = [(0.42, 1.86), (5.02, 1.86), (0.42, 3.0), (5.02, 3.0)]
    w, h = 4.42, 1.02
    for (title, items, accent), (x, y) in zip(stages, pos):
        add_rect(s, Inches(x), Inches(y), Inches(w), Inches(h), fill=WHITE, line=LINE, radius=0.1, shadow=True)
        add_rect(s, Inches(x), Inches(y), Inches(0.09), Inches(h), fill=accent, shape=MSO_SHAPE.RECTANGLE)
        add_text(s, Inches(x + 0.24), Inches(y + 0.09), Inches(w - 0.4), Inches(0.28),
                 [(title, 11.5, NAVY, True)])
        add_text(s, Inches(x + 0.24), Inches(y + 0.4), Inches(w - 0.44), Inches(0.56),
                 [[("　".join(items[:2]), 9, INK, False)], [("　".join(items[2:]), 9, INK, False)]],
                 line_spacing=1.25)
    # 双会导入条（含插图）
    band = add_rect(s, Inches(0.42), Inches(4.18), Inches(9.02), Inches(2.72),
                    fill=NAVY, radius=0.08, shadow=True)
    add_text(s, Inches(0.66), Inches(4.32), Inches(6.0), Inches(0.3),
             [("双会资源导入（摘要）", 12.5, GOLD, True)])
    pic = crop_to(img("dual_feed.png"), 3.3 / 2.24, "dual_feed_band")
    add_pic(s, pic, Inches(5.98), Inches(4.42), Inches(3.3), Inches(2.24), border=GOLD, border_w=0.75, shadow=False)
    add_text(s, Inches(0.66), Inches(4.68), Inches(5.2), Inches(1.0),
             [[("WAIC → 喂养 3F　", 10.5, GOLD, True),
               ("参展商963 · 品牌库4262 · A/B线索500", 9.5, WHITE, False)],
              [("主赛道：算力芯片 / 具身 / Agent / AIGC / 行业AI；线索抽样邀约，论坛议题→≤30人沙龙", 9, GOLD_PALE, False)]],
             line_spacing=1.35)
    add_text(s, Inches(0.66), Inches(5.78), Inches(5.2), Inches(1.0),
             [[("ChinaJoy → 喂养 5F　", 10.5, GOLD, True),
               ("可触达优先263 · 精编422", 9.5, WHITE, False)],
              [("中小游戏工作室优先；潮玩谷子 · 创作者经济；不硬追大厂总部搬迁", 9, GOLD_PALE, False)]],
             line_spacing=1.35)
    # 右栏插图
    pic = crop_to(img("rhythm_seasons.png"), 3.42 / 5.28, "rhythm_rail")
    add_pic(s, pic, Inches(9.66), Inches(0.42), Inches(3.25), Inches(5.28), )
    chip(s, Inches(9.66), Inches(5.86), Inches(3.25), Inches(0.44),
         "每月2–3场 · 全年不断档", fill=GOLD_PALE, color=GOLD_DEEP, size=10, line=None, radius=0.3)
    add_footer(s, n)


def _theme_card(s, key, x, y, w, h):
    t = THEMES[key]
    add_rect(s, Inches(x), Inches(y), Inches(w), Inches(h), fill=WHITE, line=LINE, radius=0.08, shadow=True)
    badge(s, Inches(x + 0.18), Inches(y + 0.16), Inches(0.5), key,
          fill=GOLD if key in "EF" else NAVY, size=15)
    add_text(s, Inches(x + 0.84), Inches(y + 0.12), Inches(w - 1.0), Inches(0.3),
             [[(t["name"], 13.5, NAVY, True), ("　" + t["count"] + " · " + t["span"], 10, GOLD_DEEP, True)]])
    add_text(s, Inches(x + 0.84), Inches(y + 0.44), Inches(w - 1.0), Inches(0.26),
             [(t["role"] + "　｜　" + t["res"], 9, GRAY, False)])
    # 事件芯片
    ex, ey = x + 0.18, y + 0.82
    ew = (w - 0.36 - 0.12 * (3 - 1)) / 3
    for i, (code, name, when) in enumerate(t["events"]):
        row, col = divmod(i, 3)
        cx = ex + col * (ew + 0.12)
        cy = ey + row * 0.62
        sp = add_rect(s, Inches(cx), Inches(cy), Inches(ew), Inches(0.52),
                      fill=GOLD_PALE if key in "EF" else RGBColor(0xEE, 0xF2, 0xF9),
                      radius=0.16)
        shape_text(sp, [[(code + " · " + when, 8, GOLD_DEEP, True)], [(name, 8.5, INK, False)]],
                   line_spacing=1.0)
    add_text(s, Inches(x + 0.18), Inches(y + h - 0.34), Inches(w - 0.36), Inches(0.26),
             [[("招商落点：", 9, GOLD_DEEP, True), (t["point"], 9.5, INK, False)]])


def slide_themes(prs, n, pair, part, rail_img, link_chip):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    names = "、".join([f"{k} {THEMES[k]['name']}" for k in pair])
    add_header(s, f"03 · 六大主题（{part}）", names,
               "每个主题对应双会一类资源切片，逐场口径见 P.06—P.07 日历总表")
    _theme_card(s, pair[0], 0.42, 1.86, 8.85, 2.42)
    _theme_card(s, pair[1], 0.42, 4.44, 8.85, 2.42)
    pic = crop_to(img(rail_img), 3.42 / 5.0, f"theme_{pair[0]}")
    add_pic(s, pic, Inches(9.49), Inches(0.42), Inches(3.42), Inches(5.0))
    chip(s, Inches(9.49), Inches(5.58), Inches(3.42), Inches(1.28), link_chip,
         fill=NAVY, color=GOLD_PALE, size=9.5, line=None, radius=0.12)
    add_footer(s, n)


def slide_execution(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "04 · 单场执行标准", "每场活动，同一套打法",
               "T-14 筹备 → D日执行 → T+7 回访 → 销售促成；标准化动作对应交付清单（P.14）")
    pic = crop_to(img("exec_flow.png"), 12.49 / 2.16, "exec_banner")
    add_pic(s, pic, Inches(0.42), Inches(1.82), Inches(12.49), Inches(2.16))
    steps = [
        ("T-14", "锁定嘉宾议程物料\n+ 本场资源切片", "同浦汇", "议程/嘉宾确认/报名"),
        ("T-7", "邀约确认\n+ 场地报备", "同浦汇；园区场地确认", "名单/场地单"),
        ("T-1", "彩排\n+ 看场动线", "同浦汇；园区样板间", "彩排记录"),
        ("D日", "签到→主区→看场→建群", "同浦汇执行；园区支持", "签到/意向/群"),
        ("T+1~7", "带客回访", "同浦汇主责", "回访摘要知会园区"),
        ("促成", "房源条件与成交", "园区负责销售部分", "合同/交割"),
    ]
    x = 0.42
    w = 2.0
    for i, (t, act, resp, dlv) in enumerate(steps):
        add_chevron(s, Inches(x), Inches(4.24), Inches(w), Inches(0.62),
                    [[(t, 13, WHITE if i < 4 else NAVY, True)]],
                    fill=NAVY if i < 4 else GOLD, first=(i == 0))
        card = add_rect(s, Inches(x), Inches(4.98), Inches(w), Inches(1.86),
                        fill=WHITE, line=LINE, radius=0.1, shadow=True)
        add_text(s, Inches(x + 0.1), Inches(5.08), Inches(w - 0.2), Inches(0.56),
                 [[(ln, 9, INK, True)] for ln in act.split("\n")], line_spacing=1.18)
        add_text(s, Inches(x + 0.1), Inches(5.68), Inches(w - 0.2), Inches(0.5),
                 [[("责任　", 8, GOLD_DEEP, True), (resp, 8, GRAY, False)]], line_spacing=1.15)
        add_text(s, Inches(x + 0.1), Inches(6.34), Inches(w - 0.2), Inches(0.44),
                 [[("交付　", 8, GOLD_DEEP, True), (dlv, 8, GRAY, False)]], line_spacing=1.15)
        x += w + 0.098
    add_footer(s, n)


def slide_loop(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "05 · 转化闭环（核心机制）", "到场 → 建群 → 看场 → 促成 → 复盘回流",
               "做完活动、建完群，园区在群内做销售促成；回流数据让下一场更准")
    # 闭环：上排 3 步 + 下排 3 步 + 回流箭头
    top = [("① 到场分层建群", "D日签到分层\n当场建群", NAVY),
           ("② 看场预约入群", "T+1 发看场预约\n与资料入群", NAVY),
           ("③ 带客回访", "T+1~7 同浦汇\n持续带客回访", NAVY)]
    bottom = [("⑥ 复盘回流下一场", "数据入台账\n月度核对意向", GOLD_DEEP),
              ("⑤ 意向核对与成交", "月度联席会核对\n合同/交割", NAVY_SOFT),
              ("④ 群内销售促成", "园区发布房源\n政策/条件促成交", NAVY_SOFT)]
    x = 0.42
    w = 2.96
    for i, (t, body, fill) in enumerate(top):
        add_chevron(s, Inches(x), Inches(1.92), Inches(w), Inches(1.06),
                    [[(t, 11.5, WHITE, True)]] + [[(ln, 8.5, GOLD_PALE, False)] for ln in body.split("\n")],
                    fill=fill, first=(i == 0))
        x += w + 0.11
    x = 0.42
    for i, (t, body, fill) in enumerate(bottom):
        add_chevron(s, Inches(x), Inches(4.34), Inches(w), Inches(1.06),
                    [[(t, 11.5, WHITE, True)]] + [[(ln, 8.5, GOLD_PALE, False)] for ln in body.split("\n")],
                    fill=fill, flip_h=True)
        x += w + 0.11
    # 纵向连接与回流标识
    add_rect(s, Inches(8.82), Inches(3.08), Inches(0.5), Inches(1.18),
             fill=GOLD, shape=MSO_SHAPE.DOWN_ARROW)
    add_rect(s, Inches(0.66), Inches(3.08), Inches(0.5), Inches(1.18),
             fill=GOLD, shape=MSO_SHAPE.UP_ARROW)
    add_text(s, Inches(1.2), Inches(3.5), Inches(1.0), Inches(0.3),
             [("复盘回流", 9.5, GOLD_DEEP, True)])
    center = add_rect(s, Inches(2.1), Inches(3.3), Inches(6.3), Inches(0.72),
                      fill=WHITE, line=GOLD, line_w=1.2, radius=0.5, shadow=True)
    shape_text(center, [[("同浦汇带客 · 园区销售 · 每月核对", 12.5, NAVY, True)]])
    # 底部责任边界
    bounds = [("同浦汇", "策划 · 邀约 · 带客 · 建群 · 回访摘要"),
              ("园区", "场地物业 · 样板间 · 销售成交 · 合同交割"),
              ("载体/学术", "云创基地 · 科企联/服中心 · 复旦住房中心")]
    x = 0.42
    for t, d in bounds:
        card = add_rect(s, Inches(x), Inches(5.66), Inches(2.96), Inches(1.18),
                        fill=WHITE, line=LINE, radius=0.12, shadow=True)
        add_text(s, Inches(x + 0.14), Inches(5.78), Inches(2.7), Inches(0.3),
                 [(t, 11, GOLD_DEEP, True)])
        add_text(s, Inches(x + 0.14), Inches(6.12), Inches(2.7), Inches(0.6),
                 [(d, 9, INK, False)], line_spacing=1.3)
        x += 3.07
    # 右栏插图
    pic = crop_to(img("loop_circle.png"), 3.42 / 3.42, "loop_sq")
    add_pic(s, pic, Inches(9.49), Inches(0.42), Inches(3.42), Inches(3.42))
    chip(s, Inches(9.49), Inches(4.0), Inches(3.42), Inches(0.44),
         "原则口径", fill=NAVY, color=GOLD_PALE, size=10.5, line=None, radius=0.3)
    note = add_rect(s, Inches(9.49), Inches(4.56), Inches(3.42), Inches(2.28),
                    fill=WHITE, line=LINE, radius=0.1, shadow=True)
    add_text(s, Inches(9.67), Inches(4.72), Inches(3.1), Inches(2.0),
             [[("· 我们负责带客，园区负责销售", 9.5, INK, False)],
              [("· 不承诺必带外资企业", 9.5, INK, False)],
              [("· 出海活动另议另计价", 9.5, INK, False)],
              [("· 领事到访 / 挂牌另计价，分属不同性质", 9.5, INK, False)],
              [("· 租金不做对赌与必要性要求", 9.5, INK, False)]],
             line_spacing=1.42)
    add_footer(s, n)


def slide_deliverables(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "05 · 交付清单", "交给园区的东西，可归档、可核验",
               "回访与带客由同浦汇负责；交园区为摘要知会版，不额外增加园区事务负担")
    rows = [
        ("签约后2周", "年度活动方案细化版 + 排期日历 + 资源导入台账模板", "活动专项"),
        ("每场结束后7天", "执行卡、议程、嘉宾确认、报名页、签到表、意向表、动线图、照片/通稿", "同浦汇归档"),
        ("每场结束后7天", "回访摘要（知会版）+ 建群清单", "便于园区销售跟进与群内促成"),
        ("每月", "活动数据月报 + 下月排期确认单", "联席会前"),
        ("每季 / 年末", "季度复盘、全年台账、WAIC / ChinaJoy 衔接与补库建议", "资产沉淀"),
    ]
    gf = s.shapes.add_table(len(rows) + 1, 3, Inches(0.42), Inches(1.86), Inches(8.85), Inches(3.7))
    table = style_table(gf, [1.55, 5.2, 2.1])
    table.rows[0].height = Inches(0.4)
    for c, htext in enumerate(["周期", "交付物", "说明"]):
        set_cell(table.cell(0, c), htext, size=10.5, color=GOLD_PALE, bold=True, fill=NAVY)
    for i, row in enumerate(rows):
        table.rows[i + 1].height = Inches(0.64)
        fill = WHITE if i % 2 == 0 else GOLD_PALE
        for c, val in enumerate(row):
            set_cell(table.cell(i + 1, c), val, size=9.5,
                     color=NAVY if c == 0 else INK, bold=(c == 0),
                     fill=fill, align=PP_ALIGN.LEFT if c == 1 else PP_ALIGN.CENTER)
    # 底部协同卡
    coop = add_rect(s, Inches(0.42), Inches(5.82), Inches(8.85), Inches(1.02),
                    fill=NAVY, radius=0.1, shadow=True)
    add_text(s, Inches(0.66), Inches(5.96), Inches(8.5), Inches(0.76),
             [[("请园区协同：", 10.5, GOLD, True),
               ("确认档期 · 场地物业与样板间 · 共享在谈名单 · 负责销售部分与成交交割", 10, WHITE, False)],
              [("协同单位：", 10.5, GOLD, True),
               ("云创基地主体运营 · 科企联 / 服中心载体支持 · 复旦住房中心学术支持", 10, WHITE, False)]],
             line_spacing=1.45)
    pic = crop_to(img("deliver_package.png"), 3.42 / 5.28, "deliver_rail")
    add_pic(s, pic, Inches(9.49), Inches(0.42), Inches(3.42), Inches(5.28))
    chip(s, Inches(9.49), Inches(5.86), Inches(3.42), Inches(0.44),
         "交付即沉淀 · 全年台账", fill=GOLD_PALE, color=GOLD_DEEP, size=10, line=None, radius=0.3)
    add_footer(s, n)


def slide_pricing(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "05 · 收费与付款", "费用边界清晰 · 付款唯一方案",
               "活动策划执行服务费年度打包（金额线下确认）；另计价事项不占用年度包")
    rows = [
        ("活动策划执行服务费", "年度打包（金额线下确认）", "含常规场次策划执行、物料、主持、建群与带客回访"),
        ("租金参考", "3.3 元/㎡/天", "高于周边均价，故不做租赁对赌与租赁必要性要求"),
        ("建议免租期", "1–3 个月（面议）", "降低企业决策门槛"),
        ("物业参考", "13.8 元/㎡/月", "园区统一标准"),
        ("待租房源", "8 间优先看场", "带客当场可看"),
        ("出海专题活动", "另议 · 另计价", "不在年度活动服务费内"),
        ("领事到访 / 挂牌活动", "另计价", "分属不同性质，分开计价"),
    ]
    gf = s.shapes.add_table(len(rows) + 1, 3, Inches(0.42), Inches(1.86), Inches(8.85), Inches(3.4))
    table = style_table(gf, [2.3, 2.45, 4.1])
    table.rows[0].height = Inches(0.38)
    for c, htext in enumerate(["项目", "口径", "说明"]):
        set_cell(table.cell(0, c), htext, size=10.5, color=GOLD_PALE, bold=True, fill=NAVY)
    for i, row in enumerate(rows):
        table.rows[i + 1].height = Inches(0.42)
        fill = WHITE if i % 2 == 0 else GOLD_PALE
        for c, val in enumerate(row):
            set_cell(table.cell(i + 1, c), val, size=9.5,
                     color=NAVY if c == 0 else INK, bold=(c == 0),
                     fill=fill, align=PP_ALIGN.LEFT if c == 2 else PP_ALIGN.CENTER)
    # 付款节点条
    add_text(s, Inches(0.42), Inches(5.46), Inches(6.0), Inches(0.3),
             [("付款节点（与交付验收挂钩）", 12.5, NAVY, True)])
    segs = [("签约后7日内", "50%", "启动款·锁定年度日历", 2.5, NAVY),
            ("Q3 / Q4 / 次年Q1 / 次年Q2", "各10%", "挂钩季度场次完成度与月报验收", 4.3, NAVY_SOFT),
            ("年终收官后", "10%", "尾款·年报台账影像移交", 1.9, GOLD_DEEP)]
    x = 0.42
    for t, pct, d, w, fill in segs:
        card = add_rect(s, Inches(x), Inches(5.86), Inches(w), Inches(1.0),
                        fill=fill, radius=0.12, shadow=True)
        shape_text(card, [[(pct + "　", 16, GOLD_PALE if fill != GOLD_DEEP else WHITE, True),
                           (t, 8.5, GOLD_PALE if fill != GOLD_DEEP else WHITE, False)],
                          [(d, 8, GOLD_PALE if fill != GOLD_DEEP else WHITE, False)]],
                   line_spacing=1.1)
        x += w + 0.12
    pic = crop_to(img("pricing_coins.png"), 3.42 / 5.28, "pricing_rail")
    add_pic(s, pic, Inches(9.49), Inches(0.42), Inches(3.42), Inches(5.28))
    chip(s, Inches(9.49), Inches(5.86), Inches(3.42), Inches(0.44),
         "50% + 4×10% + 10%", fill=NAVY, color=GOLD, size=11.5, line=None, radius=0.3)
    add_footer(s, n)


def slide_cases(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "06 · 往期案例", "这套打法，已经跑通过",
               "两类已交付活动实景，分别对应本年度 B/A 类场次与 F4·推广日·路演类场次")
    cases = [
        dict(title="杨「数」浦数字沙龙第七期：AI如何重塑企业DNA",
             meta="2025-06-30 · 美团上海综合指挥中心 · 政策/治理沙龙类",
             photos=["case_salon_1.jpg", "case_salon_2.jpg"],
             agenda=["参观数字化场景", "主题分享", "案例落地", "政策解读", "互动答疑"],
             orgs="杨浦区委网信办、赛博院；美团、大众点评大模型团队等",
             mapping="映射本年度：B 类政策沙龙 · A 类训练营"),
        dict(title="「融见科创·智启未来」人工智能专场路演暨投融资对接会",
             meta="2025-10 · 杨浦 · 邮储银行联合主办 · 路演对接类",
             photos=["case_roadshow_1.jpg", "case_roadshow_2.jpg"],
             agenda=["开场致辞", "科技金融致辞", "趋势分享", "6+6+6+1路演", "点评收口"],
             orgs="杨浦科创促进会、邮储银行；复楚智能、卡房信息、一造科技等",
             mapping="映射本年度：F4 Demo Day · 项目推广日 · 路演类"),
    ]
    x = 0.42
    for case in cases:
        w = 6.14
        add_rect(s, Inches(x), Inches(1.86), Inches(w), Inches(5.0), fill=WHITE, line=LINE, radius=0.06, shadow=True)
        p1 = crop_to(img(case["photos"][0]), 2.93 / 1.9, case["photos"][0][:-4])
        p2 = crop_to(img(case["photos"][1]), 2.93 / 1.9, case["photos"][1][:-4])
        add_pic(s, p1, Inches(x + 0.14), Inches(2.0), Inches(2.93), Inches(1.9), border=LINE, border_w=0.75, shadow=False)
        add_pic(s, p2, Inches(x + 3.07), Inches(2.0), Inches(2.93), Inches(1.9), border=LINE, border_w=0.75, shadow=False)
        add_text(s, Inches(x + 0.14), Inches(4.02), Inches(w - 0.28), Inches(0.56),
                 [(case["title"], 11.5, NAVY, True)], line_spacing=1.12)
        add_text(s, Inches(x + 0.14), Inches(4.58), Inches(w - 0.28), Inches(0.26),
                 [(case["meta"], 8.5, GRAY, False)])
        # 议程芯片
        ax = x + 0.14
        for a in case["agenda"]:
            chip(s, Inches(ax), Inches(4.92), Inches(1.12), Inches(0.34), a,
                 fill=GOLD_PALE, color=GOLD_DEEP, size=7.5, line=None, radius=0.35)
            ax += 1.18
        add_text(s, Inches(x + 0.14), Inches(5.4), Inches(w - 0.28), Inches(0.66),
                 [[("可见机构/企业：", 8.5, GOLD_DEEP, True), (case["orgs"], 8.5, INK, False)]],
                 line_spacing=1.25)
        mp = add_rect(s, Inches(x + 0.14), Inches(6.3), Inches(w - 0.28), Inches(0.4),
                      fill=NAVY, radius=0.3)
        shape_text(mp, [[(case["mapping"], 9, GOLD_PALE, True)]])
        x += w + 0.21
    add_footer(s, n)


def slide_conclusion(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "07 · 一页纸结论", "把一年做满，双会热度落在园区",
               "与开篇呼应：大会闭幕之后，热度有了承接地——30场排期、一套打法、一个转化闭环")
    cols = [
        ("我们交付", NAVY, ["30场按日历执行（每场≤30人）", "8月项目推广日四场点火 + WAIC承接",
                        "ChinaJoy 主题内容场次", "六大主题覆盖双会主要赛道",
                        "带客、建群、回访摘要", "月报 / 季报 / 全年台账可归档"]),
        ("请园区协同", NAVY_SOFT, ["确认档期 · 场地物业与样板间", "共享在谈名单 · 群内发布条件",
                            "负责销售部分与成交交割", "云创基地/联合会/服中心协同到位",
                            "政策礼包与券务按节奏提前到位", "月度联席会核对意向与排期"]),
        ("商务口径", GOLD_DEEP, ["付款：50% + 4×10% + 10%", "租金3.3元/㎡/天·不做对赌",
                          "建议免租期1–3个月（面议）", "出海 / 领事 / 挂牌另计价",
                          "不承诺必带外资企业", "金额线下确认·签约后2周交细化方案"]),
    ]
    x = 0.42
    for title, accent, items in cols:
        w = 2.98
        add_rect(s, Inches(x), Inches(1.86), Inches(w), Inches(4.42), fill=WHITE, line=LINE, radius=0.08, shadow=True)
        hd = add_rect(s, Inches(x), Inches(1.86), Inches(w), Inches(0.52), fill=accent, radius=0.08)
        shape_text(hd, [[(title, 13, WHITE, True)]])
        paras = [[("· " + it, 9.3, INK, False)] for it in items]
        add_text(s, Inches(x + 0.18), Inches(2.56), Inches(w - 0.36), Inches(3.6), paras, line_spacing=1.62)
        x += w + 0.115
    # 下一步行动
    nxt = add_rect(s, Inches(0.42), Inches(6.44), Inches(9.23), Inches(0.52),
                   fill=GOLD_PALE, radius=0.26)
    shape_text(nxt, [[("下一步：", 10.5, GOLD_DEEP, True),
                      ("确认档期与资源切片 → 签约 → 50%启动款 → 2周内交付年度细化方案与排期日历", 10.5, INK, True)]])
    pic = crop_to(img("conclusion_summit.png"), 3.3 / 5.28, "conclusion_rail")
    add_pic(s, pic, Inches(9.83), Inches(0.42), Inches(3.08), Inches(5.28))
    chip(s, Inches(9.83), Inches(5.86), Inches(3.08), Inches(0.44),
         "我们带客 · 园区销售", fill=NAVY, color=GOLD, size=10.5, line=None, radius=0.3)
    add_footer(s, n)


def slide_closing(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, PAGE_W, PAGE_H, fill=NAVY_DEEP, shape=MSO_SHAPE.RECTANGLE)
    s.shapes.add_picture(img("closing_skyline.png"), 0, 0, PAGE_W, PAGE_H)
    ov = add_rect(s, 0, 0, PAGE_W, Inches(4.4), fill=NAVY_DEEP, shape=MSO_SHAPE.RECTANGLE)
    set_alpha(ov, 55)
    add_text(s, Inches(0.9), Inches(0.9), Inches(11.5), Inches(0.4),
             [("THANKS · 谢谢审阅", 13, GOLD, True)])
    add_text(s, Inches(0.9), Inches(1.3), Inches(11.5), Inches(1.0),
             [("把一年做满，让双会热度落在园区", 34, WHITE, True)])
    add_rect(s, Inches(0.9), Inches(2.28), Inches(1.6), Pt(2.4), fill=GOLD, shape=MSO_SHAPE.RECTANGLE)
    add_text(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(0.4),
             [("创智汇 · 30场活动专项交付　｜　我们带客 · 园区销售 · 同步衔接 WAIC & ChinaJoy", 13, GOLD_PALE, False)])
    # 底部主体信息卡
    card = add_rect(s, Inches(0.9), Inches(5.7), Inches(11.53), Inches(1.06), fill=NAVY, radius=0.12, shadow=True)
    set_alpha(card, 82)
    add_text(s, Inches(1.14), Inches(5.86), Inches(11.1), Inches(0.76),
             [("主体运营：上海市云计算创新基地（国家级孵化器）　·　学术支持：复旦大学住房政策研究中心", 10, WHITE, False),
              ("活动及运营载体支持：杨浦区科技企业联合会 · 科技企业服务中心　·　策划执行与带客：同浦汇", 10, GOLD_PALE, False)],
             line_spacing=1.5)
    add_text(s, Inches(0.9), Inches(7.02), Inches(11.5), Inches(0.3),
             [(FOOTER + "　·　18 / 18", 8.5, GOLD_PALE, False)])


# ---------------------------------------------------------------- 主流程
def main():
    prs = Presentation()
    prs.slide_width = PAGE_W
    prs.slide_height = PAGE_H

    slide_cover(prs)                 # 01 封面
    slide_toc(prs, 2)                # 02 目录
    slide_background(prs, 3)         # 03 背景与判断
    slide_mainline(prs, 4)           # 04 总体思路
    slide_framework(prs, 5)          # 05 总体框架
    _calendar_slide(prs, 6, CAL_H1, "上", "2026年8月—12月 · 15场",
                    "calendar_h1.png",
                    ["每月2–3场不断档", "四场推广日8月点火", "E1承接WAIC余热"],
                    "F1 出海、L1 领事到访为另议另计价事项，不占年度活动包场次，不计入服务费。")
    _calendar_slide(prs, 7, CAL_H2, "下", "2026年12月中—2027年7月 · 15场",
                    "calendar_h2.png",
                    ["F4 Demo Day 年中收口", "D5 预热 WAIC 2027", "收获期CJ场次密集"],
                    "次年场次月份以“次X月”简写为“X月”标注；档期为约定期，以月度排期确认单为准。")
    slide_rhythm(prs, 8)             # 08 运营节奏
    slide_themes(prs, 9, ("A", "B"), "上", "theme_ab.png",
                 "A/B 类对应已交付案例①\n杨数浦数字沙龙 → 见 P.16")
    slide_themes(prs, 10, ("C", "D"), "中", "theme_cd.png",
                 "C/D 类对应已交付案例②\n融见科创路演大场 → 见 P.16")
    slide_themes(prs, 11, ("E", "F"), "下", "theme_ef.png",
                 "F4 为全年收口；F1/L1 另计价\n商务口径 → 见 P.15")
    slide_execution(prs, 12)         # 12 单场执行标准
    slide_loop(prs, 13)              # 13 转化闭环
    slide_deliverables(prs, 14)      # 14 交付清单
    slide_pricing(prs, 15)           # 15 收费与付款
    slide_cases(prs, 16)             # 16 往期案例
    slide_conclusion(prs, 17)        # 17 一页纸结论
    slide_closing(prs, 18)           # 18 封底

    prs.save(OUT)
    print(f"已生成: {OUT}  共 {len(prs.slides.__iter__.__self__._sldIdLst)} 页" if False else f"已生成: {OUT}")


if __name__ == "__main__":
    main()
