# -*- coding: utf-8 -*-
"""精排版白皮书简报：每页配摄影图，输出 PPTX + 单页 PNG。

版式：1920×1080。中文字体 Noto Sans CJK SC。
"""
import os
from PIL import Image, ImageDraw, ImageFont, ImageFilter, ImageEnhance

W, H = 1920, 1080
PHOTO = "/workspace/whitepaper/assets/photos"
CHART = "/workspace/whitepaper/assets/charts"
LOGO = "/workspace/whitepaper/assets/logo_fudan_hprc.png"
OUT = "/workspace/whitepaper/assets/slides"
PPTX = "/workspace/whitepaper/WAIC2026人工智能产业空间白皮书-精排版.pptx"
DL = "/workspace/whitepaper/下载版本"

NAVY = (10, 42, 82)
BLUE = (14, 78, 155)
RED = (200, 16, 46)
GOLD = (201, 162, 79)
WHITE = (255, 255, 255)
INK = (28, 33, 44)
MUTED = (92, 98, 110)
LIGHT = (246, 248, 251)
PANEL = (255, 255, 255)

FONT_R = "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"
FONT_B = "/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc"

os.makedirs(OUT, exist_ok=True)
os.makedirs(DL, exist_ok=True)


def font(path, size):
    return ImageFont.truetype(path, size)


def cover_fill(path, size=(W, H), focus="center"):
    im = Image.open(path).convert("RGB")
    tw, th = size
    scale = max(tw / im.width, th / im.height)
    im = im.resize((int(im.width * scale), int(im.height * scale)), Image.Resampling.LANCZOS)
    x = (im.width - tw) // 2
    if focus == "top":
        y = 0
    elif focus == "bottom":
        y = im.height - th
    else:
        y = (im.height - th) // 2
    return im.crop((x, y, x + tw, y + th))


def darken(im, factor=0.55):
    return ImageEnhance.Brightness(im).enhance(factor)


def gradient_bottom(base, height=420, opacity=210):
    overlay = Image.new("RGBA", base.size, (0, 0, 0, 0))
    d = ImageDraw.Draw(overlay)
    y0 = H - height
    for i in range(height):
        a = int(opacity * (i / height) ** 1.15)
        d.line([(0, y0 + i), (W, y0 + i)], fill=(8, 24, 48, a))
    out = base.convert("RGBA")
    out.alpha_composite(overlay)
    return out.convert("RGB")


def rect(draw, box, color, r=0):
    if r:
        draw.rounded_rectangle(box, radius=r, fill=color)
    else:
        draw.rectangle(box, fill=color)


def wrap(draw, text, fnt, max_w):
    lines = []
    for para in text.split("\n"):
        if not para:
            lines.append("")
            continue
        line = ""
        for ch in para:
            trial = line + ch
            if draw.textlength(trial, font=fnt) <= max_w:
                line = trial
            else:
                if line:
                    lines.append(line)
                line = ch
        lines.append(line)
    return lines


def text(draw, xy, s, fnt, fill=WHITE, anchor="lt"):
    draw.text(xy, s, font=fnt, fill=fill, anchor=anchor)


def bullets(draw, x, y, items, fnt, fill=INK, gap=46, max_w=980, bullet_color=RED):
    bf = font(FONT_B, fnt.size)
    for item in items:
        draw.ellipse([x, y + 10, x + 12, y + 22], fill=bullet_color)
        lines = wrap(draw, item, fnt, max_w)
        for i, ln in enumerate(lines):
            draw.text((x + 28, y + i * (fnt.size + 8)), ln, font=fnt, fill=fill)
        y += max(gap, len(lines) * (fnt.size + 10) + 14)
    return y


def footer(draw, page, total=22):
    f = font(FONT_R, 18)
    draw.text((56, H - 36),
              "复旦大学住房政策研究中心  ·  研究文稿第二号  FDU-HPRC-WP-2026-02",
              font=f, fill=(180, 186, 196))
    draw.text((W - 56, H - 36), f"{page:02d} / {total:02d}",
              font=f, fill=(180, 186, 196), anchor="rt")
    draw.rectangle([0, 0, 12, H], fill=GOLD)


def put_logo(canvas, xy=(48, 36), width=520, plate=False):
    logo = Image.open(LOGO).convert("RGBA")
    ratio = width / logo.width
    logo = logo.resize((width, int(logo.height * ratio)), Image.Resampling.LANCZOS)
    x, y = xy
    if plate:
        pad_x, pad_y = 18, 12
        base = canvas.convert("RGBA")
        overlay = Image.new("RGBA", canvas.size, (0, 0, 0, 0))
        od = ImageDraw.Draw(overlay)
        od.rounded_rectangle(
            [x - pad_x, y - pad_y, x + logo.width + pad_x, y + logo.height + pad_y],
            radius=16, fill=(255, 255, 255, 235))
        base.alpha_composite(overlay)
        base.paste(logo, xy, logo)
        canvas.paste(base.convert("RGB"))
        return
    canvas.paste(logo, xy, logo)


def put_chart(canvas, path, box):
    x0, y0, x1, y1 = box
    tw, th = x1 - x0, y1 - y0
    im = Image.open(path).convert("RGBA")
    scale = min(tw / im.width, th / im.height)
    im = im.resize((int(im.width * scale), int(im.height * scale)), Image.Resampling.LANCZOS)
    px = x0 + (tw - im.width) // 2
    py = y0 + (th - im.height) // 2
    canvas.paste(im, (px, py), im)


def save_slide(im, n):
    path = os.path.join(OUT, f"slide_{n:02d}.png")
    im.save(path, "PNG", optimize=True)
    print("slide", n, path)
    return path


# =====================================================================
# 版式 A：全出血封面 / 收束
# =====================================================================
def slide_cover():
    im = darken(cover_fill(f"{PHOTO}/photo_cover_lujiazui.png"), 0.62)
    im = gradient_bottom(im, 520, 230)
    d = ImageDraw.Draw(im)
    put_logo(im, (56, 48), 560, plate=True)
    rect(d, [56, 430, 220, 438], GOLD)
    text(d, (56, 470), "中心研究文稿 · 第二号", font(FONT_R, 28), WHITE)
    text(d, (56, 530), "WAIC2026", font(FONT_B, 86), WHITE)
    text(d, (56, 640), "人工智能产业空间白皮书", font(FONT_B, 64), WHITE)
    text(d, (56, 740), "AI 与产业空间融合的趋势、格局与新范式", font(FONT_R, 30), (220, 226, 234))
    text(d, (56, 980), "FDU-HPRC-WP-2026-02    二〇二六年八月 · 上海", font(FONT_R, 22), (200, 206, 214))
    return im


def slide_close():
    im = darken(cover_fill(f"{PHOTO}/photo_night_skyline.png"), 0.5)
    im = gradient_bottom(im, 700, 240)
    d = ImageDraw.Draw(im)
    put_logo(im, (56, 48), 520, plate=True)
    rect(d, [56, 430, 180, 438], GOLD)
    text(d, (56, 500), "空间将成为 AI 时代", font(FONT_B, 56), WHITE)
    text(d, (56, 580), "最诚实的计分板", font(FONT_B, 72), WHITE)
    lines = wrap(d, "哪座城市、哪个园区能让智能体安全地跑起来、让年轻人愿意留下来，产业就会在哪里生长。",
                 font(FONT_R, 28), 1200)
    y = 700
    for ln in lines:
        text(d, (56, y), ln, font(FONT_R, 28), (220, 226, 234))
        y += 42
    text(d, (56, 980), "复旦大学住房政策研究中心  ·  FDU-HPRC-WP-2026-02", font(FONT_R, 22), (200, 206, 214))
    return im


# =====================================================================
# 版式 B：左图右文
# =====================================================================
def slide_split(photo, kicker, title, items, page, extra=None):
    canvas = Image.new("RGB", (W, H), WHITE)
    left = cover_fill(photo, (760, H))
    left = darken(left, 0.82)
    canvas.paste(left, (0, 0))
    d = ImageDraw.Draw(canvas)
    rect(d, [760, 0, W, H], WHITE)
    rect(d, [760, 0, 768, H], GOLD)
    text(d, (812, 56), kicker, font(FONT_B, 22), BLUE)
    tlines = wrap(d, title, font(FONT_B, 40), 1020)
    y = 100
    for ln in tlines:
        text(d, (812, y), ln, font(FONT_B, 40), NAVY)
        y += 52
    rect(d, [812, y + 8, 912, y + 16], RED)
    y += 40
    bullets(d, 812, y, items, font(FONT_R, 26), INK, gap=52, max_w=1000)
    if extra:
        text(d, (812, 980), extra, font(FONT_R, 18), MUTED)
    footer(d, page)
    return canvas


# =====================================================================
# 版式 C：顶部摄影条 + 图表
# =====================================================================
def slide_chart(photo, kicker, title, chart, page, caption=None):
    canvas = Image.new("RGB", (W, H), LIGHT)
    banner = darken(cover_fill(photo, (W, 250), focus="top"), 0.55)
    canvas.paste(banner, (0, 0))
    d = ImageDraw.Draw(canvas)
    rect(d, [0, 242, W, 250], GOLD)
    text(d, (56, 70), kicker, font(FONT_B, 20), GOLD)
    tlines = wrap(d, title, font(FONT_B, 40), 1700)
    y = 108
    for ln in tlines:
        text(d, (56, y), ln, font(FONT_B, 40), WHITE)
        y += 50
    put_chart(canvas, chart, (40, 270, 1880, 1000))
    if caption:
        text(d, (56, 1018), caption, font(FONT_R, 16), MUTED)
    footer(d, page)
    return canvas


# =====================================================================
# 版式 D：章节扉页
# =====================================================================
def slide_chapter(photo, num, title, subtitle, page):
    im = darken(cover_fill(photo), 0.48)
    overlay = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    od = ImageDraw.Draw(overlay)
    od.rectangle([0, 0, 820, H], fill=(10, 32, 64, 200))
    im = im.convert("RGBA")
    im.alpha_composite(overlay)
    im = im.convert("RGB")
    d = ImageDraw.Draw(im)
    rect(d, [0, 0, 12, H], GOLD)
    text(d, (64, 280), num, font(FONT_B, 96), GOLD)
    tlines = wrap(d, title, font(FONT_B, 48), 680)
    y = 420
    for ln in tlines:
        text(d, (64, y), ln, font(FONT_B, 48), WHITE)
        y += 62
    slines = wrap(d, subtitle, font(FONT_R, 24), 680)
    y += 20
    for ln in slines:
        text(d, (64, y), ln, font(FONT_R, 24), (210, 216, 224))
        y += 36
    footer(d, page)
    return im


# =====================================================================
# 版式 E：四宫格卡片
# =====================================================================
def slide_cards(photo, kicker, title, cards, page):
    canvas = Image.new("RGB", (W, H), LIGHT)
    banner = darken(cover_fill(photo, (W, 220), focus="top"), 0.5)
    canvas.paste(banner, (0, 0))
    d = ImageDraw.Draw(canvas)
    text(d, (56, 58), kicker, font(FONT_B, 20), GOLD)
    text(d, (56, 108), title, font(FONT_B, 40), WHITE)
    colors = [RED, BLUE, NAVY, GOLD]
    for i, (h, body) in enumerate(cards):
        x = 48 + i * 468
        rect(d, [x, 270, x + 440, 980], WHITE, r=18)
        rect(d, [x, 270, x + 440, 286], colors[i])
        text(d, (x + 28, 320), h, font(FONT_B, 26), NAVY)
        y = 380
        for ln in wrap(d, body, font(FONT_R, 22), 380):
            text(d, (x + 28, y), ln, font(FONT_R, 22), INK)
            y += 34
    footer(d, page)
    return canvas


# =====================================================================
# 组装 22 页
# =====================================================================
slides = []

slides.append(slide_cover())  # 01

slides.append(slide_split(
    f"{PHOTO}/photo_expo_center.png",
    "编制说明  ·  FDU-HPRC-WP-2026-02",
    "一份用大会数据写成的产业空间白皮书",
    [
        "基础数据：WAIC 2026 全量资源整合总表——963 家参展商、175 场论坛、4262 家品牌库。",
        "外部交叉：xsct.ai 大模型评测榜与 watcha.cn 青年产品社区（2026 年 8 月抓取）。",
        "空间实证：本中心杨浦办公园区市场报告样本口径。",
        "文稿定位：中心研究文稿第二号，聚焦 AI 与产业空间融合。",
    ], 2))

slides.append(slide_split(
    f"{PHOTO}/photo_tech_park.png",
    "目录",
    "八章结构：从产业本体回到空间本位",
    [
        "一至三章  WAIC 全景、国内外对照、论坛风向",
        "四至五章  大模型竞争格局  ×  青年 AI 消费观察",
        "六至七章  具身智能场景  ×  2026 展望与建议",
        "第八章    AI 时代产业园区与未来办公新范式",
    ], 3, extra="图表 22 幅  ·  摄影配图每页一张  ·  精排版简报"))

slides.append(slide_split(
    f"{PHOTO}/photo_waic_crowd.png",
    "摘要  ·  五个核心判断",
    "2026：从模型竞赛转向空间落地",
    [
        "硬科技占半壁江山：机器人 24.2% + 算力 23.5%，合计 47.7%。",
        "国产大模型并跑且更便宜：综合榜 Top20 占 13 席，价差逾 4 倍。",
        "青年把 AI 驯化为生活伙伴：搭子化、情绪外包、meme 化。",
        "空间成为新的竞争变量：算力密度、场景开放度、空间柔性。",
        "中美欧三条不可互换路径：中国优势在硬件、性价比与园区。",
    ], 4))

slides.append(slide_chapter(
    f"{PHOTO}/photo_night_skyline.png",
    "01",
    "发展背景与宏观趋势",
    "治理进入多边舞台，产业进入双梯队竞合。空间，开始成为下一阶段的稀缺要素。",
    5))

slides.append(slide_chart(
    f"{PHOTO}/photo_datacenter.png",
    "国内外对照",
    "中美欧三条路径：中国优势在硬件、性价比与园区",
    f"{CHART}/chart22_cn_us_eu.png", 6,
    "示意框架，综合 WAIC 产业结构、XSCT 定价与公开治理进展"))

slides.append(slide_chart(
    f"{PHOTO}/photo_campus_towers.png",
    "国内四极分工",
    "上海均衡、北京算力、广东硬件、浙江机器人",
    f"{CHART}/chart16_regional_mix.png", 7,
    "数据来源：WAIC2026 全量资源整合总表 · 参展商注册地"))

slides.append(slide_chart(
    f"{PHOTO}/photo_humanoid_show.png",
    "产业结构",
    "963 家参展商：硬科技接近半壁江山",
    f"{CHART}/chart01_exhibitor_industry.png", 8))

slides.append(slide_chart(
    f"{PHOTO}/photo_waic_crowd.png",
    "细分赛道",
    "两个爆发点：具身智能 93 家，AI Agent 87 家",
    f"{CHART}/chart02_subsector_top15.png", 9))

slides.append(slide_chart(
    f"{PHOTO}/photo_expo_center.png",
    "空间映射",
    "展馆即产业地图：张江=底座，西岸=体验，世博=首发",
    f"{CHART}/chart15_hall_industry.png", 10))

slides.append(slide_chart(
    f"{PHOTO}/photo_office_white.png",
    "大模型格局  ·  xsct.ai",
    "综合榜 Top20：国产模型占 13 席，榜首国产",
    f"{CHART}/chart07_llm_top20.png", 11,
    "XSCT Bench 综合分 = 基础×30% + 进阶×40% + 困难×30%"))

slides.append(slide_chart(
    f"{PHOTO}/photo_datacenter.png",
    "性价比革命",
    "同样聪明，谁更便宜：价差逾 4 倍",
    f"{CHART}/chart08_llm_price_perf.png", 12,
    "国产均价约 1.43 美元/百万 token，海外约 6.27 美元"))

slides.append(slide_chart(
    f"{PHOTO}/photo_youth.png",
    "青年观察  ·  watcha.cn",
    "效率工具打底，情感与身份表达增值",
    f"{CHART}/chart12_watcha_categories.png", 13,
    "观猹热门 Top50：效率工具 18、通用助手 15；搭子化成为主流话术"))

slides.append(slide_chart(
    f"{PHOTO}/photo_humanoid_show.png",
    "具身智能  ·  H3 展厅",
    "超 150 家企业在卷场景：零售是中间地带，工业最拥挤",
    f"{CHART}/chart11_embodied_scenes.png", 14))

slides.append(slide_split(
    f"{PHOTO}/photo_service_robot.png",
    "典型场景",
    "当机器人成为空间的原住民",
    [
        "零售：银河通用便利店值守、穹彻药房取药、擎朗「具身社区」零遥操。",
        "工业：智元半导体上下料、乐聚 1:1 工厂产线、计时器取代参数表。",
        "构想：机器人友好楼宇、园区具身商业层、楼宇即训练场、酒店服务链。",
        "上海已将商业零售列入具身智能重点落地场景。",
    ], 15))

slides.append(slide_chapter(
    f"{PHOTO}/photo_office_biophilic.png",
    "08",
    "产业园区与未来办公新范式",
    "前七章回答产业发生了什么。这一章回答空间应当如何回应。",
    16))

slides.append(slide_chart(
    f"{PHOTO}/photo_office_biophilic.png",
    "园区定位重估",
    "从房东到复合运营商：六要素能力体系",
    f"{CHART}/chart14_park_radar.png", 17,
    "算力密度 · 场景开放度 · 数据要素 · 人才社群 · 资本可得性 · 空间柔性"))

slides.append(slide_chart(
    f"{PHOTO}/photo_tech_park.png",
    "杨浦实证",
    "AI 企业为大学与社群支付租金溢价——大创智最显著",
    f"{CHART}/chart20_yangpu_plates.png", 18,
    "大创智成交租金 4.88 元/㎡·天，集聚 24 家 AI/大模型企业；低租金远郊并不能自动换来 AI"))

slides.append(slide_cards(
    f"{PHOTO}/photo_campus_towers.png",
    "上海样本",
    "四种 AI 产业空间原型",
    [
        ("底座型 · 张江", "智算与芯片引擎。张江科学会堂 114 家展商中 80 家为算力芯片。产品：液冷机房、绿电套餐、超节点专线。"),
        ("体验型 · 西岸", "终端与内容消费场。西岸承办 25 场论坛，大模型展商 22 家。产品：展示层 + 消费级硬件体验。"),
        ("试验型 · 杨浦", "校地转化与空间智能。大创智 24 家 AI 企业，复兴岛承办量子城市论坛。产品：中试车间与场景开放协议。"),
        ("制度型 · 世博", "全球治理与首发经济。世博中心承载 52% 论坛及合作组织协定签署。产品：标准发布与会展转化。"),
    ], 19))

slides.append(slide_split(
    f"{PHOTO}/photo_office_white.png",
    "未来办公",
    "四种新范式：人机共生的空间语法",
    [
        "从人均工位到人机共生单元：面积按「人 + 机 + 算力」测算。",
        "从格子间到四象限：专注舱、协作场、展示层、实验区。",
        "从精装交付到柔性冗余：承重、层高、电力预留 30% 以上。",
        "从物业服务到环境即服务：物业费按服务事件而非平方米计价。",
    ], 20))

slides.append(slide_split(
    f"{PHOTO}/photo_robot_office.png",
    "人机共生",
    "给机器人留路，给模型留电，给年轻人留夜",
    [
        "给机器人留路：层高 ≥4.5m 中试层、货运电梯、充电消毒间、楼宇 5G。",
        "给模型留电：配电按训练/推理峰值预留冗余，提供可计量绿电套餐。",
        "给年轻人留夜：24 小时创作者工位、无人零售与具身社区商业。",
        "空置率高往往是产品错配，正确动作是改产品，而不是先降租。",
    ], 21))

slides.append(slide_close())  # 22

paths = []
for i, im in enumerate(slides, 1):
    paths.append(save_slide(im, i))

# ---- PPTX：每页一张全出血图，保证观感不被字体替换破坏 ----
from pptx import Presentation
from pptx.util import Inches, Emu

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]
for p in paths:
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(p, Emu(0), Emu(0), width=prs.slide_width, height=prs.slide_height)
prs.save(PPTX)
print("pptx", PPTX)
print("slides", len(paths))
