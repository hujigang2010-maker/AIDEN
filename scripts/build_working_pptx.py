"""
冠松 01# 楼 · 招商合作工作版 PPT（12 页）

给你继续深化用。收费和 90 天必须项不要改乱。
可加厚：带看照片、供应商名录、闭门会名单、市北走看。

重新生成：python3 scripts/build_working_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pathlib import Path

NAVY = RGBColor(0x1B, 0x35, 0x55)
GOLD = RGBColor(0xB7, 0x86, 0x2E)
INK = RGBColor(0x2A, 0x2D, 0x34)
CLOUD = RGBColor(0xF0, 0xF1, 0xF3)
GREEN = RGBColor(0x2F, 0x7F, 0x5B)
RED = RGBColor(0xB2, 0x3B, 0x3B)
GREY = RGBColor(0x6B, 0x73, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

CN_FONT = "WenQuanYi Micro Hei"
EN_FONT = "Georgia"
SLIDES = []


def set_run(run, text, *, size=14, bold=False, color=INK, italic=False):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = EN_FONT
    rPr = run._r.get_or_add_rPr()
    for tag in ("ea", "cs"):
        existing = rPr.find(qn(f"a:{tag}"))
        if existing is not None:
            rPr.remove(existing)
        rPr.append(rPr.makeelement(qn(f"a:{tag}"), {"typeface": CN_FONT}))


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    box.fill.background()
    box.line.fill.background()
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        set_run(run, ln, size=size, bold=bold, color=color, italic=italic)
    return box


def add_rect(slide, x, y, w, h, *, fill=NAVY):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    shp.text_frame.text = ""
    return shp


def add_round(slide, x, y, w, h, text="", *, fill=NAVY, color=WHITE,
              size=12, bold=True, align=PP_ALIGN.CENTER):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p.add_run(), text, size=size, bold=bold, color=color)
    return shp


def add_table(slide, x, y, w, h, header, rows, *, header_size=11, body_size=11,
              col_widths=None):
    ts = slide.shapes.add_table(len(rows) + 1, len(header), x, y, w, h)
    table = ts.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw
    for j, htxt in enumerate(header):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.text = ""
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p.add_run(), htxt, size=header_size, bold=True, color=WHITE)
    for i, row in enumerate(rows, start=1):
        bg = WHITE if i % 2 else CLOUD
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            cell.text = ""
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            set_run(p.add_run(), str(val), size=body_size, color=INK)
    return ts


def add_chrome(slide, prs, *, page_no, label, title, subtitle=""):
    sw, sh = prs.slide_width, prs.slide_height
    add_rect(slide, 0, 0, sw, Emu(380000), fill=NAVY)
    add_round(slide, Inches(0.45), Inches(0.18), Inches(2.5), Inches(0.36),
              label, fill=GOLD, color=NAVY, size=11)
    add_text(slide, Inches(3.15), Inches(0.10), Inches(9.5), Inches(0.52),
             title, size=22, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(3.15), Inches(0.52), Inches(9.5), Inches(0.28),
                 subtitle, size=12, color=CLOUD, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 0, sh - Emu(80000), sw, Emu(80000), fill=GOLD)
    add_text(slide, Inches(0.45), Inches(7.05), Inches(9.4), Inches(0.28),
             "工作版 v1.1 · 可深化 · 收费口径勿改乱", size=10, color=GREY)
    add_text(slide, Inches(11.3), Inches(7.05), Inches(1.6), Inches(0.28),
             f"{page_no} / 0", size=10, color=GREY, align=PP_ALIGN.RIGHT)


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    SLIDES.append(s)
    return s


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    sw, sh = prs.slide_width, prs.slide_height

    # 1 封面
    s = new_slide(prs)
    add_rect(s, 0, 0, sw, sh, fill=NAVY)
    add_rect(s, 0, Inches(5.1), sw, Emu(28000), fill=GOLD)
    add_text(s, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.4),
             "工作版 v1.1 · 给你继续改", size=16, bold=True, color=GOLD, italic=True)
    add_text(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(1.5),
             "01# 楼招商合作\n先干 90 天",
             size=40, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.5), Inches(11.7), Inches(0.7),
             "包干 8 万 · 佣金跟中介同一把尺 · 招到盖章才算完",
             size=18, color=CLOUD)
    add_text(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.8),
             "深化方向：带看照片 / 冠松供应商名录 / 闭门会名单 / 市北走看\n"
             "不要改：收费页、90 天必须项（改了和协议对不上）",
             size=14, color=CLOUD)

    # 2 听懂了
    s = new_slide(prs)
    add_chrome(s, prs, page_no=2, label="对齐",
               title="我们听懂了",
               subtitle="得到大脑：跟进招商。不是再做定位，也不是招 4S 店。")
    rows = [
        ("楼要有人租", "汽车展厅走不通。目标是生效合同，不是品牌发布会。"),
        ("不当小白鼠", "同区先看市北高新。不搞全国首创。"),
        ("不丢掉汽车", "不招整车 4S。配套研发、后市场可以留。"),
        ("先看人干活", "90 天包干。达不到必须项，可以不转正。"),
    ]
    for i, (t, b) in enumerate(rows):
        y = Inches(1.25) + Inches(1.3) * i
        add_round(s, Inches(0.5), y, Inches(0.7), Inches(1.05), str(i + 1),
                  fill=NAVY, color=WHITE, size=22)
        add_text(s, Inches(1.45), y, Inches(3.4), Inches(1.05), t,
                 size=22, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(4.9), y, Inches(7.8), Inches(1.05), b,
                 size=18, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # 3 楼硬条件（吸收实测）
    s = new_slide(prs)
    add_chrome(s, prs, page_no=3, label="楼",
               title="这栋楼能讲的硬条件",
               subtitle="吸收设计实测。带看就讲这些，不讲园区梦。可补现场照片。")
    add_table(
        s, Inches(0.4), Inches(1.2), Inches(12.5), Inches(5.4),
        ["项", "数", "招商时怎么用"],
        [
            ["用地", "C6 教育科研设计用地", "只招研发设计，不按 4S / 公寓主业态对外讲"],
            ["规模", "地上 15,152.75 ㎡ · 9F · 高 44.95 m", "考核分母约 8,300 ㎡（1F+2F 自留）"],
            ["层高", "1F 5.7 / 2F 6.3 / 3F 5.7 / 4F 5.4", "3–4F 进设备；1F+2F 展示接待，不设整车展厅"],
            ["落位", "8–9F 总部 · 3–4F 硬科技 · 6–7F 算法", "按层高和预算分，不按「链主叙事」硬塞"],
            ["绿建", "装配式 100% · 绿建二星 · 光伏 540 ㎡", "外资 / 药企 / 芯片客户的 ESG 加分"],
            ["停车", "108 个（3 普 + 105 机械）", "如实讲，不编测试场"],
        ],
        header_size=13, body_size=13,
        col_widths=[Inches(1.6), Inches(4.6), Inches(6.3)],
    )

    # 4 招谁
    s = new_slide(prs)
    add_chrome(s, prs, page_no=4, label="客群",
               title="招谁、不招谁",
               subtitle="吸收 8 案例的客群分类。名单下一页。")
    boxes = [
        (NAVY, "优先", "AI / 集成电路 / 智能机器人研发\n中心城区招人、要层高、要接待"),
        (GOLD, "可组合", "生物医药研发、医疗器械研发\n汽车配套：电控、传感、软件、测试"),
        (RED, "不招", "整车 4S / 展厅 / 卖场\n零售、教培、仓储、生产"),
    ]
    for i, (c, t, b) in enumerate(boxes):
        x = Inches(0.45) + Inches(4.2) * i
        add_rect(s, x, Inches(1.25), Inches(4.0), Inches(3.3), fill=CLOUD)
        add_rect(s, x, Inches(1.25), Inches(4.0), Inches(0.55), fill=c)
        add_text(s, x, Inches(1.25), Inches(4.0), Inches(0.55), t,
                 size=18, bold=True, color=WHITE if c != GOLD else NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.2), Inches(2.0), Inches(3.6), Inches(2.3),
                 b, size=16, color=INK, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(0.45), Inches(4.75), Inches(12.4), Inches(1.85), fill=NAVY)
    add_text(s, Inches(0.7), Inches(4.95), Inches(12.0), Inches(1.5),
             "公寓先做 C6 合规，90 天不对外承诺。\n"
             "对标只提两家：市北高新（同区已跑通）、西岸 AI 大厦（单栋冠名）。\n"
             "其余 6 个案例放包里，问到再掏。",
             size=16, color=WHITE)

    # 5 种子名单（吸收 30 家，展示 18）
    s = new_slide(prs)
    add_chrome(s, prs, page_no=5, label="名单",
               title="90 天种子客户（从案例里收）",
               subtitle="可深化：换成冠松 4S / 配件供应商名录，比这份更真。")
    cols = [
        ("AI / 芯片", NAVY, ["地平线", "黑芝麻智能", "寒武纪行歌",
                            "商汤", "MiniMax", "无问芯穹"]),
        ("汽车配套", GREEN, ["Momenta", "禾赛科技", "德赛西威",
                             "经纬恒润", "四维图新", "华域汽车电子"]),
        ("生物医药", GOLD, ["联影医疗", "微创医疗", "复宏汉霖",
                            "药明康德（分部）", "第一三共（外资）", "赛诺菲（外资）"]),
    ]
    for i, (t, c, names) in enumerate(cols):
        x = Inches(0.45) + Inches(4.2) * i
        add_rect(s, x, Inches(1.2), Inches(4.0), Inches(5.4), fill=CLOUD)
        add_rect(s, x, Inches(1.2), Inches(4.0), Inches(0.55), fill=c)
        add_text(s, x, Inches(1.2), Inches(4.0), Inches(0.55), t,
                 size=16, bold=True, color=WHITE if c != GOLD else NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        for j, n in enumerate(names):
            add_text(s, x + Inches(0.3), Inches(1.95) + Inches(0.7) * j,
                     Inches(3.4), Inches(0.6), f"{j + 1}  {n}",
                     size=16, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # 6 90天指标
    s = new_slide(prs)
    add_chrome(s, prs, page_no=6, label="90 天",
               title="90 天要交什么",
               subtitle="必须项达不到，建议不转正。锚定 TS 是冲刺，不绑死。")
    add_table(
        s, Inches(0.4), Inches(1.2), Inches(12.5), Inches(4.4),
        ["", "必须看见", "冲刺"],
        [
            ["库", "≥ 150 家，决策人到人", "200 家"],
            ["接触", "深度接触 ≥ 8 家", "12 家"],
            ["带看", "≥ 5 场，48 小时纪要", "8 场"],
            ["意向", "有效书面意向 ≥ 3 份", "意向面积 ≥ 1,500 ㎡"],
            ["渠道", "中介 2 家 + 首报制度", "第 3 家"],
            ["复盘", "第 90 天董事长 90 分钟漏斗会", "1 份锚定 Term Sheet"],
        ],
        header_size=14, body_size=14,
        col_widths=[Inches(1.8), Inches(5.6), Inches(5.1)],
    )
    add_text(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.8),
             "出 PPT、开会、只给名单，一律不算完成。",
             size=16, bold=True, color=RED)

    # 7 90天排期
    s = new_slide(prs)
    add_chrome(s, prs, page_no=7, label="排期",
               title="90 天怎么排",
               subtitle="不办 200 人发布会，不装修样板间，不跑 8 个园区。")
    phases = [
        ("D1–30 进场",
         "确认栏、钥匙图纸\n一页楼书（层高/C6）\n看房动线：1F→3–4F→8–9F\n种子库 + 中介 2 家\n周五第一份周报"),
        ("D31–60 见客",
         "种子名单攻坚\n带看 + 纪要\n15 人闭门 1 场\n只提市北、西岸\n投促办备案，不加码"),
        ("D61–90 收口",
         "意向写成书面\n绿区能定的定\n黄区 48 小时请示\n漏斗红黄灯\n转正 / 再试 / 停"),
    ]
    for i, (t, b) in enumerate(phases):
        x = Inches(0.45) + Inches(4.2) * i
        add_rect(s, x, Inches(1.25), Inches(4.0), Inches(5.35), fill=CLOUD)
        add_rect(s, x, Inches(1.25), Inches(4.0), Inches(0.6), fill=NAVY)
        add_text(s, x, Inches(1.25), Inches(4.0), Inches(0.6), t,
                 size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.2), Inches(2.05), Inches(3.6), Inches(4.3),
                 b, size=16, color=INK, align=PP_ALIGN.CENTER)

    # 8 怎么招
    s = new_slide(prs)
    add_chrome(s, prs, page_no=8, label="方法",
               title="怎么招（漏斗收干净）",
               subtitle="吸收中介首报、周报字段。90 天先开 2 家中介，不要五家一起烧。")
    steps = [
        ("1", "建库", "种子 30 家\n+ 冠松供应商"),
        ("2", "报备", "30 天首报\n一客一份"),
        ("3", "带看", "空层 + 层高尺\n不装修样板"),
        ("4", "报价", "绿区直接报\n超线必请示"),
        ("5", "书面", "意向 / TS\n用途写上 C6"),
        ("6", "盖章", "甲方盖合同\n盯保证金"),
    ]
    for i, (n, t, b) in enumerate(steps):
        x = Inches(0.4) + Inches(2.15) * i
        add_round(s, x + Inches(0.55), Inches(1.4), Inches(0.9), Inches(0.9), n,
                  fill=GOLD, color=NAVY, size=22)
        add_text(s, x, Inches(2.45), Inches(2.05), Inches(0.5), t,
                 size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(3.0), Inches(2.05), Inches(1.5), b,
                 size=13, color=INK, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(0.4), Inches(4.8), Inches(12.5), Inches(1.75), fill=CLOUD)
    add_text(s, Inches(0.65), Inches(4.95), Inches(12.1), Inches(1.5),
             "带看三句话：中心城区招人 · 3–4F 层高能进设备 · C6 合规做研发。\n"
             "不讲「中国中心城区首个智驾总部」。那是我们编过的，和他要的相反。\n"
             "每周五漏斗：新增、带看、阻塞、下周三件最重要的事。",
             size=15, color=INK)

    # 9 收费对齐
    s = new_slide(prs)
    add_chrome(s, prs, page_no=9, label="收费",
               title="月度和佣金为什么改",
               subtitle="叠在一起高了。单独看「一个点佣金」并不高。")
    add_table(
        s, Inches(0.35), Inches(1.2), Inches(12.6), Inches(3.7),
        ["", "上一版（偏高）", "现在对齐"],
        [
            ["90 天现金", "启动 15 万 + 月度 8×3 ≈ 39 万", "包干 8 万，不另收月度"],
            ["成功佣金", "首年租金 8% / 12%（听着像切利润）", "首月净租金 100% / 锚定 150%（和中介同一尺）"],
            ["中介单", "中介满佣 + 我们再拿 50%", "业主只付一套，中介约 70% + 我们约 30%"],
            ["转正后月度", "8 万/月从进场就收", "第 91 天起 4 万/月；当月有佣金则抵扣"],
        ],
        header_size=13, body_size=13,
        col_widths=[Inches(2.2), Inches(5.3), Inches(5.1)],
    )
    add_rect(s, Inches(0.35), Inches(5.1), Inches(12.6), Inches(1.5), fill=NAVY)
    add_text(s, Inches(0.55), Inches(5.25), Inches(12.2), Inches(1.25),
             "试跑 90 天若无人成交，业主只出 8 万。有成交，再按「一个点」付。\n"
             "4,000 ㎡ 空一年，6.5 元/㎡·天，大约少收 949 万。",
             size=16, color=WHITE)

    # 10 佣金举例
    s = new_slide(prs)
    add_chrome(s, prs, page_no=10, label="举例",
               title="一单大概多少钱",
               subtitle="首月不含税净租金 = 面积 × 日租金 × 30。免租不影响这个分子。")
    add_table(
        s, Inches(0.4), Inches(1.2), Inches(12.5), Inches(3.55),
        ["例子", "面积", "首月净租金", "我们自拓（100%）", "锚定（150%）"],
        [
            ["腰部研发", "800 ㎡ × 6.5 元", "约 15.6 万", "15.6 万", "—"],
            ["整层/锚定", "1,680 ㎡ × 5.8 元", "约 29.2 万", "—", "约 43.8 万"],
            ["中介带来的腰部", "800 ㎡ × 6.5 元", "约 15.6 万", "业主仍付 15.6 万\n我们约 4.7 万", "—"],
        ],
        header_size=13, body_size=13,
        col_widths=[Inches(2.4), Inches(2.4), Inches(2.3), Inches(2.7), Inches(2.7)],
    )
    add_text(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.6),
             "包干 8 万从第一笔成功佣金里抵扣。招到了，试跑费等于退回。\n"
             "甲方自己带来的未报备客户，不付成功佣金。\n"
             "可深化：按他们心里的目标租金再做一列。",
             size=15, color=INK)

    # 11 授权
    s = new_slide(prs)
    add_chrome(s, prs, page_no=11, label="授权",
               title="红黄绿 · 您不用每单上桌",
               subtitle="吸收谈判 Playbook。黄区 48 小时不批 = 不同意。")
    add_table(
        s, Inches(0.4), Inches(1.2), Inches(12.5), Inches(3.2),
        ["条款", "绿区 · 可直接报", "黄区 · 48h", "红区 · 董事长"],
        [
            ["租金 元/㎡·天", "≥ 6.5", "5.8–6.5", "底线 5.0"],
            ["免租（月）", "≤ 9", "9–15", "底线 24"],
            ["装补 元/㎡", "≤ 600", "600–1,000", "底线 1,500"],
        ],
        header_size=14, body_size=15,
        col_widths=[Inches(2.5), Inches(3.3), Inches(3.3), Inches(3.4)],
    )
    add_rect(s, Inches(0.4), Inches(4.7), Inches(12.5), Inches(1.9), fill=CLOUD)
    add_text(s, Inches(0.6), Inches(4.9), Inches(12.1), Inches(1.6),
             "合同章始终在您手里。不承诺政府政策、落户、公寓。不收租户或中介回扣。\n"
             "方案 B：先 90 天。方案 A：进场即 12 个月独家（第 91 天起月度 4 万）。\n"
             "建议勾 B。",
             size=16, color=INK)

    # 12 请定 + 深化
    s = new_slide(prs)
    add_chrome(s, prs, page_no=12, label="拍板",
               title="今天勾完，回去再深化名单",
               subtitle="数字写在 05b 确认栏。全稿随后出，不改收费结构。")
    asks = [
        ("01", "签法", "建议 B：先 90 天。"),
        ("02", "收费", "包干 8 万 · 首月 100% / 锚定 150% · 转正后月度 4 万可抵扣。"),
        ("03", "人", "唯一对接人 + 红区谁拍板 + 绿区能否直接报。"),
        ("04", "名录", "请冠松提供 4S / 配件 / 保险供应商名单，作为种子库。"),
    ]
    for i, (n, k, v) in enumerate(asks):
        y = Inches(1.25) + Inches(1.25) * i
        add_round(s, Inches(0.5), y, Inches(1.1), Inches(1.05), n,
                  fill=GOLD, color=NAVY, size=20)
        add_text(s, Inches(1.8), y, Inches(1.8), Inches(1.05), k,
                 size=20, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(3.7), y, Inches(9.0), Inches(1.05), v,
                 size=16, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    total = len(SLIDES)
    for sl in SLIDES:
        for shape in sl.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text and run.text.strip().endswith(" / 0"):
                        run.text = run.text.replace(" / 0", f" / {total}")

    out = Path(__file__).resolve().parent.parent / "docs" / "advisory" / \
        "deck" / "冠松01楼-招商合作-工作版.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"✓ Deck written: {out}  ({total} slides)")
    return out


if __name__ == "__main__":
    build()
