"""
冠松 01# 研发楼 · 招商委托说明 PPT
给董事长/决策层：为什么签委托协议、怎么替他们招、怎么算账。

重新生成：python3 scripts/build_mandate_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pathlib import Path

NAVY = RGBColor(0x1B, 0x35, 0x55)
GOLD = RGBColor(0xB7, 0x86, 0x2E)
INK = RGBColor(0x2A, 0x2D, 0x34)
CLOUD = RGBColor(0xF0, 0xF1, 0xF3)
RED = RGBColor(0xB2, 0x3B, 0x3B)
GREEN = RGBColor(0x2F, 0x7F, 0x5B)
GREY = RGBColor(0x6B, 0x73, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xCF, 0xD2, 0xD7)
STONE = RGBColor(0x8A, 0x92, 0x9C)

CN_FONT = "WenQuanYi Micro Hei"
EN_FONT = "Georgia"
SLIDES = []


def set_run(run, text, *, size=14, bold=False, color=INK, italic=False,
            font_cn=CN_FONT, font_en=EN_FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_en
    rPr = run._r.get_or_add_rPr()
    for tag in ("ea", "cs"):
        existing = rPr.find(qn(f"a:{tag}"))
        if existing is not None:
            rPr.remove(existing)
        ea = rPr.makeelement(qn(f"a:{tag}"), {"typeface": font_cn})
        rPr.append(ea)


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    box.fill.background()
    box.line.fill.background()
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        set_run(run, ln, size=size, bold=bold, color=color, italic=italic)
    return box


def add_rect(slide, x, y, w, h, *, fill=NAVY, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.5)
    shp.shadow.inherit = False
    shp.text_frame.text = ""
    return shp


def add_round(slide, x, y, w, h, text="", *, fill=NAVY, color=WHITE,
              size=12, bold=True, align=PP_ALIGN.CENTER):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p.add_run(), text, size=size, bold=bold, color=color)
    return shp


def add_table(slide, x, y, w, h, header, rows, *, header_fill=NAVY,
              header_color=WHITE, zebra=(WHITE, CLOUD), header_size=11,
              body_size=10, col_widths=None):
    cols = len(header)
    n_rows = len(rows) + 1
    ts = slide.shapes.add_table(n_rows, cols, x, y, w, h)
    table = ts.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw
    for j, htxt in enumerate(header):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Pt(4)
        tf.margin_right = Pt(4)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p.add_run(), htxt, size=header_size, bold=True, color=header_color)
    for i, row in enumerate(rows, start=1):
        bg = zebra[(i - 1) % 2]
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Pt(4)
            tf.margin_right = Pt(4)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            set_run(p.add_run(), str(val), size=body_size, color=INK)
    return ts


def add_chrome(slide, prs, *, page_no, phase_label="", page_title="", subtitle=""):
    sw, sh = prs.slide_width, prs.slide_height
    add_rect(slide, 0, 0, sw, Emu(380000), fill=NAVY)
    if phase_label:
        add_round(slide, Inches(0.5), Inches(0.18), Inches(2.6), Inches(0.36),
                  phase_label, fill=GOLD, color=NAVY, size=11, bold=True)
    if page_title:
        add_text(slide, Inches(3.3), Inches(0.10), Inches(9.3), Inches(0.55),
                 page_title, size=22, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(3.3), Inches(0.55), Inches(9.3), Inches(0.28),
                 subtitle, size=11, color=CLOUD, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 0, sh - Emu(80000), sw, Emu(80000), fill=GOLD)
    add_text(slide, Inches(0.5), Inches(7.05), Inches(9.0), Inches(0.30),
             "冠松 · 01# 研发楼 · 招商服务委托说明 v1.0（保密）",
             size=9, color=GREY)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(2.0), Inches(0.30),
             f"{page_no} / 0", size=9, color=GREY, align=PP_ALIGN.RIGHT)


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    SLIDES.append(s)
    return s


def card(slide, x, y, w, h, title, body, *, title_fill=NAVY):
    add_rect(slide, x, y, w, h, fill=CLOUD)
    add_rect(slide, x, y, w, Emu(90000), fill=title_fill)
    add_text(slide, x, y, w, Inches(0.36), title,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x + Inches(0.12), y + Inches(0.42), w - Inches(0.24),
             h - Inches(0.5), body, size=12, color=INK)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    sw, sh = prs.slide_width, prs.slide_height

    # 1 封面
    s = new_slide(prs)
    add_rect(s, 0, 0, sw, sh, fill=NAVY)
    add_rect(s, 0, Inches(5.05), sw, Emu(30000), fill=GOLD)
    deco = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.0), Inches(-2.5),
                              Inches(8), Inches(8))
    deco.fill.solid()
    deco.fill.fore_color.rgb = RGBColor(0x25, 0x40, 0x64)
    deco.line.fill.background()
    deco.shadow.inherit = False

    add_text(s, Inches(0.8), Inches(0.7), Inches(11), Inches(0.4),
             "GS · iDrive Hub · Mandate", size=14, bold=True,
             color=GOLD, italic=True)
    add_text(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.5),
             "我们替你们招",
             size=48, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.15), Inches(11.5), Inches(0.7),
             "01# 研发楼 · 招商服务委托协议说明",
             size=24, color=GOLD)
    add_text(s, Inches(0.8), Inches(4.0), Inches(11.5), Inches(0.55),
             "不是顾问备忘录 · 不是中介居间 · 是独家全案操盘",
             size=16, color=CLOUD)

    add_round(s, Inches(0.8), Inches(5.5), Inches(3.2), Inches(0.45),
              "给董事长拍板 v1.0", fill=GOLD, color=NAVY, size=13, bold=True)
    add_text(s, Inches(4.2), Inches(5.5), Inches(8), Inches(0.45),
             "配套：docs/legal/05-合作协议-招商服务委托协议.docx",
             size=13, color=CLOUD, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.8), Inches(6.4), Inches(11), Inches(0.4),
             "费用为建议口径 · 签约前法务定稿 · 完成标准只有一条：生效租赁合同",
             size=12, color=STONE, italic=True)

    # 2 判断
    s = new_slide(prs)
    add_chrome(s, prs, page_no=2, phase_label="判断",
               page_title="不要再签顾问协议",
               subtitle="方案已经有了。楼要填满，差的是谁在授权里把客户带到盖章。")

    add_table(
        s, Inches(0.5), Inches(1.25), Inches(12.3), Inches(3.4),
        ["", "顾问协议（兜圈子）", "本委托协议（直接干）"],
        [
            ["完成标准", "出方案、出案例、开会", "生效租赁合同（盖章+保证金到账）"],
            ["乙方角色", "参谋", "替业主招商的操盘方"],
            ["对入驻率", "不负责", "M6 / M12 节点写进合同"],
            ["收费", "咨询费，与空置无关", "月度干活 + 招到才抽佣"],
            ["报价", "每次请示", "绿区可直接报，黄区 48 小时，红区董事长"],
            ["业态", "建议，可听可不听", "禁止整车展厅，写进合同"],
        ],
        header_size=13, body_size=12,
        col_widths=[Inches(2.0), Inches(5.15), Inches(5.15)],
    )

    add_round(s, Inches(0.5), Inches(4.9), Inches(12.3), Inches(1.7),
              "一句话：顾问解决「招什么」；委托解决「谁去招、招到哪算完、怎么付钱」。\n"
              "01–04 号协议是园区对外的。05 号才是冠松委托我们操盘的这一份。",
              fill=CLOUD, color=NAVY, size=16, bold=False, align=PP_ALIGN.LEFT)

    # 3 一页纸
    s = new_slide(prs)
    add_chrome(s, prs, page_no=3, phase_label="骨架",
               page_title="协议四句话",
               subtitle="能在会上念完，就不要再加形容词")

    items = [
        ("01 谁", "甲方冠松 = 业主。乙方 = 独家招商操盘方。\n租赁合同仍由甲方盖章，乙方不能私自签。"),
        ("02 干什么", "建库、攻坚、带看、谈判、统管中介、协同政府、盯到保证金到账。\n出 PPT 不算完成。"),
        ("03 算完成", "生效合同。M6：1 份锚定 TS 或意向 ≥ 2,000 ㎡。\nM12：签约 ≥ 3,735 ㎡ 或入驻率 ≥ 45%。"),
        ("04 怎么付钱", "启动费进场 · 月度费干活 · 成功佣金招到才收。\n中介佣金另付，乙方不截留。"),
    ]
    for i, (k, v) in enumerate(items):
        x = Inches(0.5) + Inches(6.15) * (i % 2)
        y = Inches(1.3) + Inches(2.6) * (i // 2)
        add_rect(s, x, y, Inches(5.95), Inches(2.4), fill=CLOUD)
        add_round(s, x + Inches(0.2), y + Inches(0.2), Inches(2.2), Inches(0.45),
                  k, fill=NAVY, color=WHITE, size=14)
        add_text(s, x + Inches(0.2), y + Inches(0.8), Inches(5.55), Inches(1.4),
                 v, size=15, color=INK)

    # 4 干什么
    s = new_slide(prs)
    add_chrome(s, prs, page_no=4, phase_label="动作",
               page_title="替你们招 · 八件事",
               subtitle="每一件都要留下甲方能抽查的记录")

    actions = [
        ("1 建库", "≥ 200 家目标企业\n决策人到人"),
        ("2 攻坚", "建联、拜访、痛点\n接触记录可核验"),
        ("3 带看", "现场讲解 + 48h 纪要\n异议点写下来"),
        ("4 谈判", "绿区直接报\nTS / 合同商务稿"),
        ("5 中介", "统一报备、冲突裁决\n只保留一套真相"),
        ("6 政府", "材料、陪同、诉求清单\n对外承诺须甲方确认"),
        ("7 活动", "闭门会 / 发布会招商段\n硬成本甲方承担"),
        ("8 盯签约", "从意向到盖章\n到保证金到账"),
    ]
    for i, (t, b) in enumerate(actions):
        x = Inches(0.45) + Inches(3.15) * (i % 4)
        y = Inches(1.3) + Inches(2.7) * (i // 4)
        add_rect(s, x, y, Inches(3.0), Inches(2.5), fill=CLOUD)
        add_rect(s, x, y, Inches(3.0), Inches(0.55), fill=NAVY)
        add_text(s, x, y, Inches(3.0), Inches(0.55), t,
                 size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.12), y + Inches(0.7), Inches(2.76), Inches(1.65),
                 b, size=14, color=INK, align=PP_ALIGN.CENTER)

    # 5 KPI
    s = new_slide(prs)
    add_chrome(s, prs, page_no=5, phase_label="考核",
               page_title="12 个月 · 招到什么算完成",
               subtitle="分母：可出租净面积约 8,300 ㎡（1F+2F 大堂展厅自留，不计入）")

    add_table(
        s, Inches(0.45), Inches(1.25), Inches(12.4), Inches(3.55),
        ["时点", "过程（必须看得见）", "结果（合同认的）"],
        [
            ["M3", "库 ≥ 200 · 深接触 ≥ 8 · 带看 ≥ 5 · 周报不断档", "≥ 3 份书面意向"],
            ["M6", "中介报备制度跑通 · 至少 1 场闭门招商", "1 份锚定 TS，或意向面积 ≥ 2,000 ㎡"],
            ["M9", "漏斗周会连续 · 重点客户一页纸齐备", "生效签约 ≥ 3,000 ㎡（约 36%）"],
            ["M12", "档案与 SOP 可移交", "生效签约 ≥ 3,735 ㎡，或入驻率 ≥ 45%"],
        ],
        header_size=13, body_size=13,
        col_widths=[Inches(1.3), Inches(6.3), Inches(4.8)],
    )

    add_rect(s, Inches(0.45), Inches(5.05), Inches(12.4), Inches(1.55), fill=CLOUD)
    add_text(s, Inches(0.65), Inches(5.15), Inches(12.0), Inches(1.35),
             "M6 杀招（保护业主）：达不到结果指标，且不是甲方不配合导致的，甲方可终止，之后月度费不付。\n"
             "M12 未达目标不自动算根本违约，但影响续约；续约则第 13–18 月月度费下浮 20%。\n"
             "只开会、只出方案、只有口头意向 —— 一律不算完成。",
             size=14, color=INK)

    # 6 授权 + 业态
    s = new_slide(prs)
    add_chrome(s, prs, page_no=6, phase_label="授权",
               page_title="红黄绿 + 业态边界",
               subtitle="董事长不用每单上桌；也不能被口头加码")

    add_table(
        s, Inches(0.4), Inches(1.2), Inches(12.5), Inches(2.7),
        ["条款", "绿区 · 可直接报", "黄区 · 48h 批复", "红区 · 董事长"],
        [
            ["租金 元/㎡·天", "≥ 6.5", "5.8–6.5", "底线 5.0"],
            ["免租（月）", "≤ 9", "9–15", "底线 24"],
            ["装补 元/㎡", "≤ 600", "600–1,000", "底线 1,500"],
        ],
        header_size=12, body_size=13,
        col_widths=[Inches(2.4), Inches(3.3), Inches(3.4), Inches(3.4)],
    )

    add_text(s, Inches(0.5), Inches(4.1), Inches(6.0), Inches(0.4),
             "允许招", size=14, bold=True, color=GREEN)
    add_text(s, Inches(0.5), Inches(4.5), Inches(6.0), Inches(1.8),
             "AI / 集成电路 / 智能机器人研发\n"
             "生物医药与器械研发（依法定许可）\n"
             "汽车配套研发（电控、传感、软件、测试）\n"
             "非整车销售",
             size=14, color=INK)

    add_text(s, Inches(7.0), Inches(4.1), Inches(5.8), Inches(0.4),
             "禁止招", size=14, bold=True, color=RED)
    add_text(s, Inches(7.0), Inches(4.5), Inches(5.8), Inches(1.8),
             "整车 4S / 整车展厅 / 汽车卖场\n"
             "纯零售、教培、仓储、生产\n"
             "地上主业态公寓（人才短租须先做 C6 合规）\n"
             "回扣、口头政策、包牌照、保 GMV",
             size=14, color=INK)

    # 7 三笔账
    s = new_slide(prs)
    add_chrome(s, prs, page_no=7, phase_label="算账",
               page_title="三笔账 · 招到才抽佣",
               subtitle="建议口径，写入附件 E 可改数字，不改结构")

    fees = [
        ("启动费 · 进场", "15 万一次性", "10–20 万可谈\n全额抵扣成功佣金\n不是再收一笔咨询费", NAVY),
        ("月度费 · 干活", "8 万 / 月", "6–12 万可谈\n3–4 人核心组\n低于自建团队全成本", NAVY),
        ("成功佣金 · 结果", "8% / 链主 12%", "基数 = 首年租金\n不含税、物业，扣免租\n没签成不抽", GOLD),
    ]
    for i, (t, n, b, c) in enumerate(fees):
        x = Inches(0.45) + Inches(4.2) * i
        add_rect(s, x, Inches(1.25), Inches(4.0), Inches(3.35), fill=CLOUD)
        add_rect(s, x, Inches(1.25), Inches(4.0), Inches(0.5), fill=c)
        add_text(s, x, Inches(1.25), Inches(4.0), Inches(0.5), t,
                 size=14, bold=True, color=WHITE if c == NAVY else NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, Inches(1.85), Inches(4.0), Inches(0.7), n,
                 size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.2), Inches(2.6), Inches(3.6), Inches(1.8),
                 b, size=14, color=INK, align=PP_ALIGN.CENTER)

    add_rect(s, Inches(0.45), Inches(4.8), Inches(12.4), Inches(1.8), fill=NAVY)
    add_text(s, Inches(0.65), Inches(4.95), Inches(12.0), Inches(1.5),
             "和空置比：4,000 ㎡ 空一年，按 6.5 元/㎡·天约损失 949 万。\n"
             "付给操盘方：月度 96 万 + 成功佣金约 57 万 ≈ 153 万（启动费抵扣进佣金）。\n"
             "中介佣金甲方另付。甲方自有未报备客户成交，不付成功佣金。",
             size=15, color=WHITE)

    # 8 分工
    s = new_slide(prs)
    add_chrome(s, prs, page_no=8, phase_label="分工",
               page_title="甲方拍板，乙方跑动",
               subtitle="独家 12 个月，但租赁合同章永远在甲方手里")

    add_rect(s, Inches(0.45), Inches(1.25), Inches(6.05), Inches(5.35), fill=CLOUD)
    add_rect(s, Inches(0.45), Inches(1.25), Inches(6.05), Inches(0.55), fill=NAVY)
    add_text(s, Inches(0.45), Inches(1.25), Inches(6.05), Inches(0.55),
             "甲方 · 冠松", size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.7), Inches(2.0), Inches(5.6), Inches(4.3),
             "产权与出租主体\n"
             "指定唯一对接人，黄区 48h 答复\n"
             "红区由董事长或授权人批\n"
             "提供看房、图纸、物业陪同\n"
             "链主拜访按预约出席\n"
             "法务 10 个工作日反馈合同\n"
             "不向已报备客户暗中压价\n"
             "中介佣金、活动硬成本自担",
             size=15, color=INK)

    add_rect(s, Inches(6.8), Inches(1.25), Inches(6.05), Inches(5.35), fill=CLOUD)
    add_rect(s, Inches(6.8), Inches(1.25), Inches(6.05), Inches(0.55), fill=GOLD)
    add_text(s, Inches(6.8), Inches(1.25), Inches(6.05), Inches(0.55),
             "乙方 · 操盘方", size=18, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(7.05), Inches(2.0), Inches(5.6), Inches(4.3),
             "项目经理每周现场 ≥ 3 天\n"
             "漏斗、带看、谈判、周报月报\n"
             "绿区报价，超线必请示\n"
             "统管中介报备，避免一客多报\n"
             "起草 TS / 合同商务条款（含 C6）\n"
             "不向租户、中介收回扣\n"
             "不口头承诺政府政策\n"
             "终止后 10 日移交全部档案",
             size=15, color=INK)

    # 9 保护
    s = new_slide(prs)
    add_chrome(s, prs, page_no=9, phase_label="保护",
               page_title="业主五道闸",
               subtitle="独家不是把楼交给别人瞎承诺")

    gates = [
        ("1", "M6 杀招", "达不到锚定 TS 或 2,000 ㎡ 意向，可终止，之后月度费停。"),
        ("2", "合同甲方盖章", "乙方无权以甲方名义签租赁合同、收款。"),
        ("3", "黄区默认否", "48 小时不批复 = 不同意，防止沉默即授权。"),
        ("4", "红线可解除", "虚假报备、收回扣、超授权承诺政策，立即解除并追偿。"),
        ("5", "业态写死", "禁止整车展厅；公寓不是地上主业态，须先做 C6 合规。"),
    ]
    for i, (n, t, b) in enumerate(gates):
        y = Inches(1.25) + Inches(1.05) * i
        add_round(s, Inches(0.5), y, Inches(0.7), Inches(0.85), n,
                  fill=NAVY, color=WHITE, size=20)
        add_text(s, Inches(1.4), y, Inches(3.2), Inches(0.85), t,
                 size=18, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(4.6), y, Inches(8.2), Inches(0.85), b,
                 size=16, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # 10 拍板
    s = new_slide(prs)
    add_chrome(s, prs, page_no=10, phase_label="拍板",
               page_title="今天请定四件事",
               subtitle="法务只改表述和主体名称；商务结构不要改回顾问")

    asks = [
        ("01", "路径", "走委托招商，不走顾问续约。"),
        ("02", "数字", "确认附件 E：启动费 15 万 / 月度 8 万 / 佣金 8%·12%。"),
        ("03", "人", "指定唯一对接人 + 红区最终决策人。"),
        ("04", "权", "书面确认绿区可直接报价，否则每单开会、节奏会断。"),
    ]
    for i, (n, k, v) in enumerate(asks):
        y = Inches(1.3) + Inches(1.25) * i
        add_round(s, Inches(0.55), y, Inches(1.1), Inches(1.0), n,
                  fill=GOLD, color=NAVY, size=20)
        add_text(s, Inches(1.9), y, Inches(1.8), Inches(1.0), k,
                 size=20, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(3.8), y, Inches(8.8), Inches(1.0), v,
                 size=18, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # 11 收尾
    s = new_slide(prs)
    add_rect(s, 0, 0, sw, sh, fill=NAVY)
    add_rect(s, 0, Inches(3.55), sw, Emu(30000), fill=GOLD)
    add_text(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.5),
             "一句话", size=14, bold=True, color=GOLD, italic=True,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), Inches(2.1), Inches(11.7), Inches(1.1),
             "我们替你们招。招到才算完。",
             size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), Inches(4.0), Inches(11.7), Inches(1.4),
             "协议：docs/legal/05-合作协议-招商服务委托协议.docx\n"
             "要点：docs/advisory/09-招商委托协议要点.md",
             size=16, color=CLOUD, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.5),
             "草案 v1.0 · 冠松 01# 研发楼 · 签约前法务定稿",
             size=14, color=GOLD, align=PP_ALIGN.CENTER, italic=True)

    total = len(SLIDES)
    for sl in SLIDES:
        for shape in sl.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text and run.text.strip().endswith(" / 0"):
                        run.text = run.text.replace(" / 0", f" / {total}")

    out = Path(__file__).resolve().parent.parent / "docs" / "advisory" / \
        "deck" / "冠松01楼-招商委托说明.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"✓ Deck written: {out}  ({total} slides)")
    return out


if __name__ == "__main__":
    build()
