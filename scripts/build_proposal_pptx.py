#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""生成给同浦汇的业务承接策划案 PPT（16:9，16 页）。"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))
import proposal_data as D

NAVY = RGBColor(0x0E, 0x22, 0x40)
NAVY2 = RGBColor(0x1B, 0x3A, 0x6B)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
GOLD_PALE = RGBColor(0xF4, 0xEB, 0xD3)
CREAM = RGBColor(0xFA, 0xF6, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x24, 0x30, 0x44)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LINE = RGBColor(0xE7, 0xDD, 0xC6)
TEAL = RGBColor(0x1F, 0x6B, 0x5C)
RED = RGBColor(0x8B, 0x2E, 0x2E)

FONT = "微软雅黑"
PAGE_W = Inches(13.333)
PAGE_H = Inches(7.5)
TOTAL = 16
FOOTER = "杨浦区科技企业服务中心 × 同浦汇 · 业务承接策划案 · 2026.09"

OUT = Path(__file__).resolve().parents[1] / "output" / "同浦汇_30场活动与科技企业服务中心筹备_业务承接策划案.pptx"


def _ea(run, font=FONT):
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", font)


def style_run(run, size, color, bold=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    _ea(run)


def add_text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, size=14, color=INK, bold=False, spacing=1.05):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    first = True
    for item in lines:
        if isinstance(item, tuple):
            text, sz, col, b = item
        else:
            text, sz, col, b = item, size, color, bold
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(2)
        p.line_spacing = spacing
        run = p.add_run()
        run.text = text
        style_run(run, sz, col, b)
    return box


def rect(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background() if line is None else setattr(sh.line, "color", None)
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def round_rect(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.9)
    sh.adjustments[0] = 0.08
    sh.shadow.inherit = False
    return sh


def bar(slide, x, y, w, h, fill=GOLD):
    return rect(slide, x, y, w, h, fill)


def footer(slide, page):
    rect(slide, Inches(0), Inches(7.28), PAGE_W, Inches(0.22), NAVY)
    add_text(slide, Inches(0.4), Inches(7.28), Inches(11.2), Inches(0.22),
             [(FOOTER, 9, WHITE, False)], anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(11.6), Inches(7.28), Inches(1.4), Inches(0.22),
             [(f"{page:02d} / {TOTAL:02d}", 9, GOLD, True)], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def header(slide, kicker, title, subtitle=None):
    rect(slide, Inches(0), Inches(0), PAGE_W, Inches(0.08), GOLD)
    add_text(slide, Inches(0.45), Inches(0.18), Inches(12.4), Inches(0.28),
             [(kicker, 11, GOLD, True)])
    add_text(slide, Inches(0.45), Inches(0.42), Inches(12.4), Inches(0.42),
             [(title, 24, NAVY, True)])
    if subtitle:
        add_text(slide, Inches(0.45), Inches(0.86), Inches(12.4), Inches(0.32),
                 [(subtitle, 13, GRAY, False)])
    bar(slide, Inches(0.45), Inches(1.18) if subtitle else Inches(0.92), Inches(1.4), Inches(0.045))


def blank():
    prs = blank.prs
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, Inches(0), Inches(0), PAGE_W, PAGE_H, CREAM)
    return slide


def card(slide, x, y, w, h, title, body, title_color=NAVY, accent=GOLD):
    round_rect(slide, x, y, w, h, WHITE, LINE)
    rect(slide, x, y, Inches(0.08), h, accent)
    add_text(slide, x + Inches(0.22), y + Inches(0.12), w - Inches(0.32), Inches(0.32),
             [(title, 13, title_color, True)])
    if isinstance(body, str):
        body = [body]
    add_text(slide, x + Inches(0.22), y + Inches(0.44), w - Inches(0.32), h - Inches(0.52),
             [(t, 11, INK, False) for t in body], size=11, spacing=1.12)


def kpi_chip(slide, x, y, w, h, num, label):
    round_rect(slide, x, y, w, h, WHITE, LINE)
    add_text(slide, x, y + Inches(0.12), w, Inches(0.42),
             [(num, 20, GOLD, True)], align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.08), y + Inches(0.52), w - Inches(0.16), Inches(0.36),
             [(label, 11, GRAY, False)], align=PP_ALIGN.CENTER)


def build():
    prs = Presentation()
    prs.slide_width = PAGE_W
    prs.slide_height = PAGE_H
    blank.prs = prs

    # 01 封面
    s = blank()
    rect(s, Inches(0), Inches(0), PAGE_W, PAGE_H, NAVY)
    rect(s, Inches(0), Inches(0), Inches(0.18), PAGE_H, GOLD)
    add_text(s, Inches(0.7), Inches(1.15), Inches(11.5), Inches(0.35),
             [("提交对象　" + D.DOC_FOR, 14, GOLD, False)])
    add_text(s, Inches(0.7), Inches(1.7), Inches(12), Inches(1.3),
             [(D.DOC_TITLE, 32, WHITE, True),
              (D.DOC_SUBTITLE, 28, GOLD, True)])
    add_text(s, Inches(0.7), Inches(3.35), Inches(11.5), Inches(0.7),
             [("承接范围：30 场可核验活动　＋　科技企业服务中心挂牌筹备", 16, WHITE, False),
              ("新赛道：" + D.NEW_POSITIONING + "　｜　政策包装：" + D.POLICY_PACKAGING, 14, GOLD_PALE, False)])
    chips = [
        ("30 场", "年度活动包"),
        ("30 万", "对园区打包价"),
        ("≤30 人", "单场核验"),
        ("90 天", "服中心筹备"),
        ("70 / 30", "执行 / 接口"),
    ]
    for i, (n, l) in enumerate(chips):
        x = Inches(0.7) + Inches(i * 2.4)
        round_rect(s, x, Inches(4.4), Inches(2.2), Inches(1.05), NAVY2)
        add_text(s, x, Inches(4.48), Inches(2.2), Inches(0.5),
                 [(n, 20, GOLD, True)], align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(4.95), Inches(2.2), Inches(0.4),
                 [(l, 11, WHITE, False)], align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(5.75), Inches(11.5), Inches(0.9),
             [(f"{D.DOC_FROM}　｜　{D.DOC_COFROM}", 13, WHITE, False),
              (f"学术支持：{D.DOC_SUPPORT.split('（')[0]}　｜　{D.DOC_DATE}　｜　周期 {D.DOC_PERIOD}", 12, GOLD_PALE, False),
              (D.CONFIDENTIAL, 11, RGBColor(0x9A, 0xA3, 0xB2), False)])
    add_text(s, Inches(11.6), Inches(7.05), Inches(1.4), Inches(0.3),
             [("01 / 16", 10, GOLD, True)], align=PP_ALIGN.RIGHT)

    # 02 昨日共识
    s = blank()
    header(s, "01 　背景", "昨天沟通后，需要立刻承接的两件事",
           "2026-08-31　创智汇赛道调整交流（约 36 分钟）· 合同口径已变、运营成本已发生")
    left = D.YESTERDAY[:6]
    right = D.YESTERDAY[6:]
    round_rect(s, Inches(0.4), Inches(1.35), Inches(6.15), Inches(5.7), WHITE, LINE)
    add_text(s, Inches(0.6), Inches(1.48), Inches(5.8), Inches(0.35),
             [("赛道为什么改", 14, GOLD, True)])
    add_text(s, Inches(0.6), Inches(1.88), Inches(5.8), Inches(4.95),
             [(f"{i+1}.  {t}", 12, INK, False) for i, t in enumerate(left)], spacing=1.18)
    round_rect(s, Inches(6.75), Inches(1.35), Inches(6.15), Inches(5.7), WHITE, LINE)
    add_text(s, Inches(6.95), Inches(1.48), Inches(5.8), Inches(0.35),
             [("新方向怎么落地", 14, GOLD, True)])
    add_text(s, Inches(6.95), Inches(1.88), Inches(5.8), Inches(4.95),
             [(f"{i+7}.  {t}", 12, INK, False) for i, t in enumerate(right)], spacing=1.18)
    footer(s, 2)

    # 03 承接范围
    s = blank()
    header(s, "02 　范围", "我们承接同浦汇的两项具体业务", "不是替换同浦汇的园区入口地位，而是把执行和筹备接下来")
    for i, (t, b) in enumerate(D.TAKEOVER):
        x = Inches(0.4) + Inches((i % 4) * 3.2)
        y = Inches(1.4)
        round_rect(s, x, y, Inches(3.05), Inches(2.15), WHITE, LINE)
        rect(s, x, y, Inches(3.05), Inches(0.08), GOLD)
        add_text(s, x + Inches(0.15), y + Inches(0.25), Inches(2.75), Inches(0.55),
                 [(f"0{i+1}  {t}", 14, NAVY, True)])
        add_text(s, x + Inches(0.15), y + Inches(0.85), Inches(2.75), Inches(1.15),
                 [(b, 12, INK, False)])
    add_text(s, Inches(0.45), Inches(3.7), Inches(12), Inches(0.35),
             [("明确不承接 / 不承诺", 14, RED, True)])
    for i, t in enumerate(D.NOT_TAKEOVER):
        col, row = i % 2, i // 2
        x = Inches(0.4) + Inches(col * 6.45)
        y = Inches(4.1) + Inches(row * 0.5)
        round_rect(s, x, y, Inches(6.3), Inches(0.44), GOLD_PALE)
        add_text(s, x + Inches(0.15), y, Inches(6.0), Inches(0.44),
                 [(t, 12, INK, False)], anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 3)

    # 04 分工
    s = blank()
    header(s, "03 　角色", "同浦汇管关系，服中心管交付", "园区管销售促成　｜　业主管场地条件　｜　复旦提供学术支持")
    colors = [GOLD, TEAL, NAVY2, NAVY, GRAY, INK]
    for i, ((who, what), c) in enumerate(zip(D.ROLES, colors)):
        y = Inches(1.38) + Inches(i * 0.9)
        round_rect(s, Inches(0.45), y, Inches(12.4), Inches(0.82), WHITE, LINE)
        rect(s, Inches(0.45), y, Inches(0.12), Inches(0.82), c)
        add_text(s, Inches(0.8), y + Inches(0.08), Inches(3.2), Inches(0.66),
                 [(who, 16, NAVY, True)], anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(4.1), y + Inches(0.08), Inches(8.5), Inches(0.66),
                 [(what, 13, INK, False)], anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 4)

    # 05 总览
    s = blank()
    header(s, "04 　产品", "两条工作包，一张图看完", "对园区报价不变　｜　对同浦汇把执行责任写清楚")
    round_rect(s, Inches(0.4), Inches(1.4), Inches(6.15), Inches(5.55), WHITE, LINE)
    rect(s, Inches(0.4), Inches(1.4), Inches(6.15), Inches(0.7), NAVY)
    add_text(s, Inches(0.6), Inches(1.5), Inches(5.8), Inches(0.5),
             [("工作包 A　30 场活动全案", 18, WHITE, True)])
    items_a = [
        "周期：2026.08—2027.07　·　约每月 2–3 场",
        "打包：30 场 / 30 万元（门票赞助可冲抵成本）",
        "核验：每场 ≤30 人，负责人 ≤30%，全年约 600+",
        "主线：智能建造能力 + 出海准备（国内可交付场）",
        "付款：签约 7 日内 50% + 四季度各 10% + 年终 10%",
        "不含：出海专题、领事到访、挂牌仪式",
        "转化：带到场、建群；销售促成仍在园区",
    ]
    add_text(s, Inches(0.65), Inches(2.25), Inches(5.7), Inches(4.4),
             [(t, 14, INK, False) for t in items_a], spacing=1.28)
    round_rect(s, Inches(6.75), Inches(1.4), Inches(6.15), Inches(5.55), WHITE, LINE)
    rect(s, Inches(6.75), Inches(1.4), Inches(6.15), Inches(0.7), TEAL)
    add_text(s, Inches(6.95), Inches(1.5), Inches(5.8), Inches(0.5),
             [("工作包 B　服中心筹备", 18, WHITE, True)])
    items_b = [
        "90 天完成挂牌材料、制度、2 人配置",
        "载体认定 + 创新券服务机构两项前置",
        "高新辅导 SOP（企业自行申报，不承诺获批）",
        "500㎡ 建筑出海展陈大纲 + 零碳研学模型",
        "联合复旦白皮书开题与授牌预热",
        "筹备不另向同浦汇收费",
        "政策收益建议：同浦汇 38% / 服中心 62%",
    ]
    add_text(s, Inches(7.0), Inches(2.25), Inches(5.7), Inches(4.4),
             [(t, 14, INK, False) for t in items_b], spacing=1.28)
    footer(s, 5)

    # 06 主题结构
    s = blank()
    header(s, "05 　活动结构", "30 场按六条线重排，服务新赛道",
           "原双会 AI 内容线不再作为主叙事；科技能力保留为智能建造的工具")
    for i, (name, n, when, what) in enumerate(D.THEMES):
        col, row = i % 3, i // 3
        x = Inches(0.4) + Inches(col * 4.25)
        y = Inches(1.4) + Inches(row * 2.7)
        round_rect(s, x, y, Inches(4.05), Inches(2.5), WHITE, LINE)
        add_text(s, x + Inches(0.2), y + Inches(0.2), Inches(3.65), Inches(0.4),
                 [(name, 16, NAVY, True)])
        add_text(s, x + Inches(0.2), y + Inches(0.65), Inches(3.65), Inches(0.4),
                 [(f"{n} 场　·　{when}", 13, GOLD, True)])
        add_text(s, x + Inches(0.2), y + Inches(1.15), Inches(3.65), Inches(1.1),
                 [(what, 13, INK, False)])
    footer(s, 6)

    # 07 日历上
    s = blank()
    header(s, "06 　排期（上）", "2026.08—2026.12　共 16 场", "8 月必须点火：赛道说明、渠道、金融、政府、进博承接")
    first = [e for e in D.EVENTS if e[1].startswith("2026")]
    # table-like rows
    y0 = Inches(1.32)
    headers = ["编号", "月份", "线条", "活动名称", "本场作用"]
    widths = [0.9, 1.15, 1.4, 5.3, 3.5]
    x0 = 0.4
    rect(s, Inches(x0), y0, Inches(sum(widths)), Inches(0.32), NAVY)
    x = x0
    for h, w in zip(headers, widths):
        add_text(s, Inches(x), y0, Inches(w), Inches(0.32),
                 [(h, 11, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
        x += w
    for i, (code, month, theme, name, why) in enumerate(first):
        y = y0 + Inches(0.32) + Inches(i * 0.325)
        bg = WHITE if i % 2 == 0 else GOLD_PALE
        rect(s, Inches(x0), y, Inches(sum(widths)), Inches(0.325), bg)
        vals = [code, month, theme, name, why]
        x = x0
        for v, w in zip(vals, widths):
            add_text(s, Inches(x + 0.04), y, Inches(w - 0.04), Inches(0.325),
                     [(v, 10, INK, False)], anchor=MSO_ANCHOR.MIDDLE)
            x += w
    footer(s, 7)

    # 08 日历下
    s = blank()
    header(s, "07 　排期（下）", "2027.01—2027.07　共 14 场", "上半年深耕能力与政策，下半年收口转化并预热下一年度")
    second = [e for e in D.EVENTS if e[1].startswith("2027")]
    y0 = Inches(1.32)
    rect(s, Inches(x0), y0, Inches(sum(widths)), Inches(0.32), NAVY)
    x = x0
    for h, w in zip(headers, widths):
        add_text(s, Inches(x), y0, Inches(w), Inches(0.32),
                 [(h, 11, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
        x += w
    for i, (code, month, theme, name, why) in enumerate(second):
        y = y0 + Inches(0.32) + Inches(i * 0.36)
        bg = WHITE if i % 2 == 0 else GOLD_PALE
        rect(s, Inches(x0), y, Inches(sum(widths)), Inches(0.36), bg)
        vals = [code, month, theme, name, why]
        x = x0
        for v, w in zip(vals, widths):
            add_text(s, Inches(x + 0.04), y, Inches(w - 0.04), Inches(0.36),
                     [(v, 11, INK, False)], anchor=MSO_ANCHOR.MIDDLE)
            x += w
    footer(s, 8)

    # 09 单场 + 转化
    s = blank()
    header(s, "08 　打法", "单场可复制，转化不压到同浦汇一个人身上",
           D.EVENT_STANDARD["节奏"])
    std = [
        ("人数核验", D.EVENT_STANDARD["人数"]),
        ("负责人", D.EVENT_STANDARD["负责人占比"]),
        ("全年触达", D.EVENT_STANDARD["触达"]),
        ("交付物", D.EVENT_STANDARD["交付"]),
    ]
    for i, (t, b) in enumerate(std):
        x = Inches(0.4) + Inches(i * 3.2)
        kpi_chip(s, x, Inches(1.4), Inches(3.05), Inches(1.15), t, b)
    steps = ["到场", "建群", "看场", "同浦汇回访", "园区促成", "复盘回流"]
    add_text(s, Inches(0.45), Inches(2.75), Inches(12), Inches(0.35),
             [("转化闭环（活动只对带到场负责，不对租金去化对赌）", 14, NAVY, True)])
    for i, st in enumerate(steps):
        x = Inches(0.4) + Inches(i * 2.15)
        round_rect(s, x, Inches(3.2), Inches(1.95), Inches(0.7), NAVY if i % 2 == 0 else TEAL)
        add_text(s, x, Inches(3.2), Inches(1.95), Inches(0.7),
                 [(f"{i+1}  {st}", 13, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < len(steps) - 1:
            add_text(s, x + Inches(1.85), Inches(3.3), Inches(0.35), Inches(0.5),
                     [("→", 16, GOLD, True)])
    add_text(s, Inches(0.45), Inches(4.15), Inches(12.4), Inches(2.7), [
        ("另计价，不进 30 万包", 14, GOLD, True),
        ("　·　".join([f"{a}（{c}）" for a, _, c in D.EXTRA_PRICE]), 13, INK, False),
        ("", 8, INK, False),
        ("同浦汇在闭环里的位置", 14, GOLD, True),
        ("客户跟踪与回访台账由同浦汇多承担；服中心重心放在下一场策划与现场。园区拿回访摘要即可，不要求逐户回访。", 13, INK, False),
    ])
    footer(s, 9)

    # 10 服中心 90 天
    s = blank()
    header(s, "09 　服中心", "90 天把科技企业服务中心筹备成可挂牌、可申报",
           "挂牌仪式另计价　｜　筹备过程不另向同浦汇收费")
    for i, (when, title, body) in enumerate(D.CENTER_90):
        x = Inches(0.4) + Inches(i * 3.2)
        round_rect(s, x, Inches(1.4), Inches(3.05), Inches(3.55), WHITE, LINE)
        rect(s, x, Inches(1.4), Inches(3.05), Inches(0.9), NAVY if i < 3 else TEAL)
        add_text(s, x + Inches(0.15), Inches(1.48), Inches(2.75), Inches(0.35),
                 [(when, 12, GOLD, True)])
        add_text(s, x + Inches(0.15), Inches(1.82), Inches(2.75), Inches(0.4),
                 [(title, 16, WHITE, True)])
        add_text(s, x + Inches(0.18), Inches(2.5), Inches(2.7), Inches(2.25),
                 [(body, 13, INK, False)])
    add_text(s, Inches(0.45), Inches(5.15), Inches(12.4), Inches(1.85), [
        ("政策上限（测算，不是保底）", 14, NAVY, True),
        f"载体年上限 {D.POLICY_CAP['载体年上限']}　｜　活动年上限 {D.POLICY_CAP['活动年上限']}",
        f"十年三部分合计上限 {D.POLICY_CAP['十年三部分上限']}",
        f"前置条件：{D.POLICY_CAP['前置条件']}",
    ])
    footer(s, 10)

    # 11 空间
    s = blank()
    header(s, "10 　空间", "6600㎡ 还是创智汇，产品要换成新赛道",
           f"{D.SPACE['区位']}　｜　办公 {D.SPACE['办公租金']}　｜　物业 {D.SPACE['物业']}")
    blocks = [
        ("3F　" + D.SPACE["3F"], "智能建造办公 / 装备体验 / 辅导接待 / OPC 工位", NAVY),
        ("5F　" + D.SPACE["5F"], "建筑出海展陈 / 模块化样品 / 新材料 / 培训沙龙", TEAL),
        ("专题展区　" + D.SPACE["展厅专题"], "昨日会上明确：先把 500㎡ 做成可参观、可讲解、可研学的主题馆", GOLD),
    ]
    for i, (t, b, c) in enumerate(blocks):
        y = Inches(1.4) + Inches(i * 1.15)
        round_rect(s, Inches(0.45), y, Inches(12.4), Inches(1.05), WHITE, LINE)
        rect(s, Inches(0.45), y, Inches(0.14), Inches(1.05), c)
        add_text(s, Inches(0.85), y + Inches(0.12), Inches(11.7), Inches(0.4),
                 [(t, 16, NAVY, True)])
        add_text(s, Inches(0.85), y + Inches(0.52), Inches(11.7), Inches(0.42),
                 [(b, 13, INK, False)])
    add_text(s, Inches(0.5), Inches(4.95), Inches(12.3), Inches(2.0), [
        ("与招商的关系（写进承接函，避免审核方误解）", 14, GOLD, True),
        "租金 3.3 元高于周边，故本方案不做租赁对赌、不做租赁必要性。活动的任务是把对的人带到场、建群、看场。",
        "待租约 8 间由园区销售闭环；佣金按 2 个月净租金（首年不重复）。建议园区给 1–3 个月免租，线下谈。",
        "零碳体验馆可做成面向中小学的研学产品，形成不完全依赖补贴的运营入口。",
    ])
    footer(s, 11)

    # 12 政策
    s = blank()
    header(s, "11 　政策", "新赛道仍然走科技企业服务，不另开一套「地产补贴」故事",
           "会上共识：智能施工、绿色节能写入经营范围后，走高新技术企业认定")
    cards = [
        ("对企业", ["先改经营范围：智能施工 / 绿色节能 / 数字化设计", "企业自行申报高新，服中心提供辅导 SOP", "不承诺获批，成功费以到账为准"]),
        ("对载体", ["争取认定为成果转化服务平台 / 孵化基地", "服中心申请成为创新券服务机构", "YOUNG立方活动按投入 50%、年封顶 200 万申报"]),
        ("对品牌", ["联合复旦发布智能建造 / 出海白皮书", "战略合作签约 + 官方授牌（仪式另计价）", "进博、领馆作为加购，不写进年包承诺"]),
    ]
    for i, (t, lines) in enumerate(cards):
        x = Inches(0.4) + Inches(i * 4.25)
        round_rect(s, x, Inches(1.4), Inches(4.05), Inches(3.4), WHITE, LINE)
        rect(s, x, Inches(1.4), Inches(4.05), Inches(0.6), NAVY)
        add_text(s, x, Inches(1.48), Inches(4.05), Inches(0.45),
                 [(t, 16, WHITE, True)], align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.2), Inches(2.2), Inches(3.65), Inches(2.4),
                 [(f"·  {ln}", 14, INK, False) for ln in lines], spacing=1.25)
    add_text(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.95), [
        ("必须写进给审核方材料的三句话", 14, RED, True),
        "1. 我们不是在招「付不起租金的施工队来租办公室」。",
        "2. 我们是在用科技企业服务中心，把智能建造产品、模块化建筑和绿色低碳能力组织成可出海的集群。",
        "3. 业主是杨浦科创集团；政府对接用科技与产业服务口径，不用「中建」名义。",
    ])
    footer(s, 12)

    # 13 商务
    s = blank()
    header(s, "12 　商务", "对园区的报价一口价，付款只保留一套", "方案二已取消　｜　与季度场次完成度、月报验收挂钩")
    rows = [
        ("活动年包", D.COMMERCIAL["活动年包"]),
        ("付款节点", D.COMMERCIAL["付款"]),
        ("验收挂钩", D.COMMERCIAL["挂钩"]),
        ("门票 / 赞助", D.COMMERCIAL["门票赞助"]),
        ("招商佣金", D.COMMERCIAL["佣金"]),
        ("服中心筹备", D.COMMERCIAL["筹备费"]),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(1.35) + Inches(i * 0.72)
        round_rect(s, Inches(0.45), y, Inches(12.4), Inches(0.64), WHITE, LINE)
        add_text(s, Inches(0.7), y, Inches(2.4), Inches(0.64),
                 [(k, 14, GOLD, True)], anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(3.2), y, Inches(9.4), Inches(0.64),
                 [(v, 13, INK, False)], anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 13)

    # 14 分成
    s = blank()
    header(s, "13 　结算", "同浦汇保留接口价值，服中心按执行结算", "两套分成互不混用：活动年包 ≠ 政策收益")
    round_rect(s, Inches(0.4), Inches(1.4), Inches(6.15), Inches(3.6), WHITE, LINE)
    rect(s, Inches(0.4), Inches(1.4), Inches(6.15), Inches(0.6), NAVY)
    add_text(s, Inches(0.4), Inches(1.48), Inches(6.15), Inches(0.45),
             [("活动年包内部结算（建议）", 16, WHITE, True)], align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(2.2), Inches(5.6), Inches(2.5), [
        ("服中心　70%", 28, GOLD, True),
        ("策划、邀约协同、现场、月报、可核验交付", 13, INK, False),
        ("", 8, INK, False),
        ("同浦汇　30%", 22, NAVY, True),
        ("园区接口、客户跟踪、回访台账、线索转交", 13, INK, False),
    ])
    round_rect(s, Inches(6.75), Inches(1.4), Inches(6.15), Inches(3.6), WHITE, LINE)
    rect(s, Inches(6.75), Inches(1.4), Inches(6.15), Inches(0.6), TEAL)
    add_text(s, Inches(6.75), Inches(1.48), Inches(6.15), Inches(0.45),
             [("政策收益分成（锁版）", 16, WHITE, True)], align=PP_ALIGN.CENTER)
    add_text(s, Inches(7.05), Inches(2.2), Inches(5.6), Inches(2.5), [
        ("同浦汇　38%  ｜  服中心　62%", 22, GOLD, True),
        ("载体申报乙方偏高；活动类可均分；申报服务费以执行方为主。", 13, INK, False),
        ("十年上限约 6920 万是可申请测算，不是保证到账。", 13, INK, False),
        ("成功费以资金到账为准，不向企业承诺。", 13, INK, False),
    ])
    add_text(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.75), [
        ("给同浦汇的一句话", 14, GOLD, True),
        "您继续做创智汇面前的招服入口和客户主人；30 场能不能办完、服中心能不能挂上牌，由服中心按本方案承接交付。",
        "这样您可以把精力放在回访和园区销售协同上，也有一份能在下周交给审核方的执行答卷。",
    ])
    footer(s, 14)

    # 15 90天给中建包
    s = blank()
    header(s, "14 　过审", "下周就能带去审核方工程部门的材料包", "业主口径：杨浦科创集团　｜　不写中建为合同业主")
    pack = [
        ("01", "一页纸定位", "智能建造 × 建筑产业出海，科技赋能，不是地产招租"),
        ("02", "30 场总表", "本 PPT 第 7–8 页 + Excel 执行台账"),
        ("03", "服中心 90 天", "挂牌、制度、载体、创新券四件套时间表"),
        ("04", "商务一页", "30 万、付款 50+4×10+10、不做租金对赌"),
        ("05", "边界清单", "出海 / 领事 / 挂牌另计价；不承诺外企"),
        ("06", "复旦协同函", "白皮书开题意向 + 授牌预热节奏（正式签约另走）"),
    ]
    for i, (n, t, b) in enumerate(pack):
        col, row = i % 3, i // 3
        x = Inches(0.4) + Inches(col * 4.25)
        y = Inches(1.4) + Inches(row * 2.55)
        round_rect(s, x, y, Inches(4.05), Inches(2.35), WHITE, LINE)
        add_text(s, x + Inches(0.2), y + Inches(0.2), Inches(3.65), Inches(0.4),
                 [(n, 18, GOLD, True)])
        add_text(s, x + Inches(0.2), y + Inches(0.65), Inches(3.65), Inches(0.45),
                 [(t, 16, NAVY, True)])
        add_text(s, x + Inches(0.2), y + Inches(1.15), Inches(3.65), Inches(0.95),
                 [(b, 13, INK, False)])
    footer(s, 15)

    # 16 下一步
    s = blank()
    header(s, "15 　请同浦汇确认", "五件事确认后，7 日内可启动", "启动条件：范围确认 + 结算确认 + 档期确认 + 口径确认 + 启动款")
    for i, t in enumerate(D.NEXT_STEPS):
        y = Inches(1.4) + Inches(i * 0.72)
        round_rect(s, Inches(0.45), y, Inches(12.4), Inches(0.64), WHITE, LINE)
        oval = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.65), y + Inches(0.14), Inches(0.36), Inches(0.36))
        oval.fill.solid()
        oval.fill.fore_color.rgb = GOLD
        oval.line.fill.background()
        add_text(s, Inches(0.65), y + Inches(0.14), Inches(0.36), Inches(0.36),
                 [(str(i + 1), 12, NAVY, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.2), y, Inches(11.4), Inches(0.64),
                 [(t, 14, INK, False)], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.5), Inches(5.15), Inches(12.3), Inches(1.8), [
        (f"联系人（同浦汇）：{D.PARTIES['联系人']}", 13, NAVY, True),
        f"我方：{D.DOC_FROM}　｜　{D.DOC_COFROM}",
        f"周期 {D.DOC_PERIOD}　｜　{D.DOC_DATE}",
        D.CONFIDENTIAL,
    ])
    footer(s, 16)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"PPT 已写入 {OUT}　页数 {len(prs.slides)}")
    return OUT


if __name__ == "__main__":
    build()
