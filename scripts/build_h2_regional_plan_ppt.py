"""绿城华东区 · 2026下半年 AI 商业化落地合作方案 PPT"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
from pathlib import Path

GREEN = RGBColor(0x00, 0x6B, 0x3F)
DARK_GREEN = RGBColor(0x00, 0x4A, 0x2C)
LIGHT_GREEN = RGBColor(0xD8, 0xE8, 0xDF)
GOLD = RGBColor(0xC8, 0xA2, 0x5B)
DARK = RGBColor(0x1F, 0x1F, 0x1F)
GREY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GREY = RGBColor(0xF5, 0xF7, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xEE, 0xF4, 0xF0)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL = 15
OUT = Path("/workspace/deliverables/绿城华东区-2026下半年AI商业化落地合作方案.pptx")


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="微软雅黑"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
        rPr = r._r.get_or_add_rPr()
        ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", font)
        cs = etree.SubElement(rPr, qn("a:cs"))
        cs.set("typeface", font)
    return tb


def add_title_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.85), GREEN)
    add_rect(slide, 0, Inches(0.85), SLIDE_W, Inches(0.05), GOLD)
    add_text(slide, Inches(0.45), Inches(0.12), Inches(12), Inches(0.6),
             title, size=22, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.45), Inches(0.95), Inches(12.2), Inches(0.35),
                 subtitle, size=11, color=GREY)


def add_footer(slide, page_no):
    add_text(slide, Inches(0.45), Inches(7.08), Inches(10), Inches(0.28),
             "绿城华东区 × 人工智能商业化落地  |  2026下半年片区战略合作方案（内部洽谈稿）",
             size=9, color=GREY)
    add_text(slide, Inches(11.4), Inches(7.08), Inches(1.5), Inches(0.28),
             f"{page_no:02d} / {TOTAL:02d}", size=9, color=GREY, align=PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, title, body, title_size=14, body_size=11):
    add_rect(slide, x, y, w, h, LIGHT_GREY)
    add_rect(slide, x, y, Inches(0.08), h, GOLD)
    add_text(slide, x + Inches(0.2), y + Inches(0.12), w - Inches(0.3), Inches(0.35),
             title, size=title_size, bold=True, color=DARK_GREEN)
    add_text(slide, x + Inches(0.2), y + Inches(0.5), w - Inches(0.3), h - Inches(0.6),
             body, size=body_size, color=DARK)


def add_table(slide, x, y, w, h, headers, rows, col_widths_ratio=None,
              header_fill=GREEN, body_size=10, header_size=11):
    n_cols = len(headers)
    n_rows = len(rows) + 1
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    tbl = tbl_shape.table
    if col_widths_ratio:
        total = sum(col_widths_ratio)
        for i, r in enumerate(col_widths_ratio):
            tbl.columns[i].width = int(w * r / total)

    def fill_cell(cell, text, fill, color, size, bold=False):
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
        cell.margin_left = Inches(0.06)
        cell.margin_right = Inches(0.06)
        cell.margin_top = Inches(0.03)
        cell.margin_bottom = Inches(0.03)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "微软雅黑"
        rPr = r._r.get_or_add_rPr()
        ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", "微软雅黑")

    for c, hdr in enumerate(headers):
        fill_cell(tbl.cell(0, c), hdr, header_fill, WHITE, header_size, True)
    for r_i, row in enumerate(rows):
        bg = WHITE if r_i % 2 == 0 else SOFT
        for c, val in enumerate(row):
            fill_cell(tbl.cell(r_i + 1, c), val, bg, DARK, body_size)
    return tbl_shape


def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, DARK_GREEN)
    add_rect(s, 0, Inches(5.9), SLIDE_W, Inches(1.6), GREEN)
    add_rect(s, Inches(0.7), Inches(1.55), Inches(1.2), Inches(0.08), GOLD)
    add_text(s, Inches(0.7), Inches(1.8), Inches(11.5), Inches(0.5),
             "绿城华东区战略合作洽谈稿", size=16, color=GOLD, bold=True)
    add_text(s, Inches(0.7), Inches(2.35), Inches(11.8), Inches(1.2),
             "2026下半年人工智能商业化落地\n与绿城片区项目合作方案",
             size=32, bold=True, color=WHITE)
    add_text(s, Inches(0.7), Inches(4.0), Inches(11.5), Inches(0.9),
             "聚焦黄浦 · 虹口 · 青浦\n以北外滩一滴水为枢纽，链接AI大厂、TMT圈层与高净值客户\n帮助区域总形成片区共识、客群资产与行政价值网络",
             size=15, color=LIGHT_GREEN)
    add_text(s, Inches(0.7), Inches(6.2), Inches(11.5), Inches(0.9),
             "主办协同：上海市杨浦区科技企业联合会  |  见微知海新质商业生态\n交付形式：PPT方案 + Excel权益与排期表  |  用途：华东区总经理一对一沟通",
             size=12, color=WHITE)


def slide_one_pager(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_title_bar(s, "01  给华东区总的一页纸：为什么值得做",
                  "不是单项目活动赞助，而是帮您经营黄浦—虹口—青浦片区的客户与行政链接")
    cards = [
        ("对您的核心价值",
         "把AI圈层、大厂高管、基金与TMT决策人，转化为绿城片区可复用的客群资产与政商互动接口。"),
        ("为什么是区域总视角",
         "活动不绑死单一案场，而以片区轮动+统一品牌叙事，强化华东区对多项目的统筹话语权与共识。"),
        ("我们怎么帮您落地",
         "一滴水主场办行/办大众 + 杨浦虹口大厂参访 + 项目开放日，形成“看见—对话—成交”闭环。"),
        ("建议先拍板什么",
         "先定档位与主推片区（建议虹口北外滩枢纽+黄浦/青浦轮动），再锁下半年4–6场关键节点。"),
    ]
    for i, (t, b) in enumerate(cards):
        x = Inches(0.45) + (i % 2) * Inches(6.35)
        y = Inches(1.5) + (i // 2) * Inches(2.5)
        add_card(s, x, y, Inches(6.1), Inches(2.25), t, b, 15, 13)
    add_footer(s, 2)


def slide_insight(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_title_bar(s, "02  区域洞察：AI 正在改写片区经营逻辑",
                  "华东区总需要的不是“再办一场会”，而是可沉淀的圈层与交易入口")
    points = [
        ("客群迁移", "高净值与企业客户越来越通过产业社群、闭门沙龙、大厂参访获取信任，而非传统案场到访。"),
        ("行政协同", "区级经信/科委/商会与园区平台，偏好有产业主题的跨项目活动，区域总出面最合适。"),
        ("项目协同", "黄浦城市更新、虹口北外滩、青浦新城等叙事不同，需要统一“绿城华东×AI人居”母品牌串起来。"),
        ("交易机会", "AI大厂、TMT、基金既是潜在买家/租户，也可能带来企业采购、联合办公、产业导入。"),
    ]
    for i, (t, b) in enumerate(points):
        y = Inches(1.45) + i * Inches(1.25)
        add_rect(s, Inches(0.45), y, Inches(12.4), Inches(1.1), LIGHT_GREY)
        add_rect(s, Inches(0.45), y, Inches(0.12), Inches(1.1), GREEN)
        add_text(s, Inches(0.8), y + Inches(0.15), Inches(2.2), Inches(0.8),
                 t, size=16, bold=True, color=DARK_GREEN, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(3.2), y + Inches(0.2), Inches(9.3), Inches(0.75),
                 b, size=13, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, 3)


def slide_engines(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_title_bar(s, "03  战略定位：三大价值引擎",
                  "客户链接 × 行政链接 × 交易链接 —— 全部服务华东区片区经营")
    engines = [
        ("① 客户链接", "GOLD",
         "把AI/TMT高管、创业者、基金LP/GP、企业主变成绿城可触达客群\n\n"
         "· 一滴水闭门宴 / 精品沙龙\n· 项目夜访与私董会\n· 高净值家庭日与企业包场"),
        ("② 行政链接", "GREEN",
         "借活动帮助区域总与区级平台、协会、园区建立常态互动\n\n"
         "· 区级领导/协会出席接口\n· 产业主题座谈会\n· 片区共识发布与媒体叙事"),
        ("③ 交易链接", "DARK",
         "不止曝光，更推动看房、企业落位、合作洽谈与后续商机\n\n"
         "· 大厂/TMT定向参访\n· 供需对接会与路演\n· 成交线索沉淀与回访"),
    ]
    colors = [GOLD, GREEN, DARK_GREEN]
    for i, ((t, _, b), c) in enumerate(zip(engines, colors)):
        x = Inches(0.4) + i * Inches(4.25)
        add_rect(s, x, Inches(1.5), Inches(4.05), Inches(5.1), LIGHT_GREY)
        add_rect(s, x, Inches(1.5), Inches(4.05), Inches(0.7), c)
        add_text(s, x + Inches(0.2), Inches(1.6), Inches(3.65), Inches(0.5),
                 t, size=18, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.25), Inches(2.45), Inches(3.55), Inches(3.9),
                 b, size=13, color=DARK)
    add_footer(s, 4)


def slide_map(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_title_bar(s, "04  片区共识地图：黄浦 · 虹口 · 青浦",
                  "以虹口北外滩为一滴水主场，黄浦与青浦做主题轮动，形成华东区统一叙事")
    rows = [
        ["虹口 / 北外滩", "主场枢纽", "一滴水系列峰会、晚宴、闭门沙龙",
         "潮鸣外滩等项目承接参观与私宴", "强化北外滩高端人居与圈层地标"],
        ["黄浦", "城市更新/企业客群", "CBD企业主沙龙、AI×商业空间专题",
         "黄浦项目开放日 / 城市更新对话", "链接中心城区决策人与企业总部"],
        ["青浦", "新城产业导入", "产业园区×AI应用落地交流",
         "青浦项目企业包场与家庭日", "服务新城置业与企业落位叙事"],
    ]
    add_table(s, Inches(0.4), Inches(1.45), Inches(12.5), Inches(3.4),
              ["片区", "角色", "主打活动", "项目联动", "对区域总的意义"],
              rows, [1.4, 1.5, 3.2, 3.0, 3.4], body_size=11, header_size=12)
    add_card(s, Inches(0.4), Inches(5.1), Inches(12.5), Inches(1.55),
             "片区共识打法（建议向总经理强调）",
             "每一场活动都同时露出“绿城华东区”母品牌 + 当期主场项目；季度末做一次三区联动复盘会，把线索、嘉宾、行政互动沉淀为华东区资产，而不是项目部各自为战。",
             14, 13)
    add_footer(s, 5)


def slide_venue(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_title_bar(s, "05  主场枢纽：北外滩一滴水",
                  "距离近、调性高、适合办行也适合办大众，天然服务虹口及周边片区")
    items = [
        ("办“行”", "闭门私董会、AI商业化圆桌、GP/LP晚宴\n人数精、议题深、转化强"),
        ("办“大众”", "公开主题峰会、媒体日、青年科技人才交流\n扩声量、扩圈层、做品牌"),
        ("办“高质量”", "大厂高管对话、产业应用案例路演\n保证嘉宾级别与内容密度"),
        ("办“转化”", "会后即赴项目参观 / 企业对接\n从场内热度变成案场动作"),
    ]
    for i, (t, b) in enumerate(items):
        x = Inches(0.4) + (i % 4) * Inches(3.2)
        add_rect(s, x, Inches(1.55), Inches(3.05), Inches(3.6), LIGHT_GREY)
        add_rect(s, x, Inches(1.55), Inches(3.05), Inches(0.65), GREEN if i % 2 == 0 else GOLD)
        add_text(s, x + Inches(0.15), Inches(1.65), Inches(2.75), Inches(0.45),
                 t, size=16, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.15), Inches(2.45), Inches(2.75), Inches(2.4),
                 b, size=13, color=DARK)
    add_text(s, Inches(0.45), Inches(5.4), Inches(12.4), Inches(1.2),
             "建议节奏：每月/每双月固定一滴水节点，黄浦与青浦各安排1–2场主题外溢，形成“主场稳定 + 片区轮动”。\n"
             "已有合作基础：2026上半年潮鸣外滩晚宴冠名已验证绿城与AI峰会的现场联动模式，下半年可升级为片区体系。",
             size=13, color=DARK)
    add_footer(s, 6)


def slide_matrix(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_title_bar(s, "06  下半年产品矩阵：四类高质量活动",
                  "可单独赞助，也可打包成华东区年度合作")
    rows = [
        ["A. 一滴水旗舰峰会/晚宴", "虹口北外滩", "AI商业化+人居/投资破局", "品牌高光+高端客群", "1–2场"],
        ["B. 精品闭门沙龙", "一滴水/项目会所", "垂直议题（算力/应用/出海等）", "深度关系与线索", "2–3场"],
        ["C. 大厂 & TMT 参访", "杨浦·虹口产业带", "腾讯云/字节/商汤/AI硬件等", "行政与产业链接", "2–3场"],
        ["D. 片区项目联动日", "黄浦/虹口/青浦", "参观+座谈+晚宴", "直接服务去化与落位", "3场轮动"],
    ]
    add_table(s, Inches(0.35), Inches(1.45), Inches(12.6), Inches(3.6),
              ["产品线", "主阵地", "内容方向", "对绿城价值", "建议频次"],
              rows, [2.6, 1.8, 3.2, 2.6, 1.4], body_size=12, header_size=12)
    add_card(s, Inches(0.35), Inches(5.25), Inches(12.6), Inches(1.4),
             "摇人能力（我们侧承诺边界）",
             "依托杨浦科技企业联合会、见微知海生态与AI峰会嘉宾池，优先邀请人工智能大厂、TMT高管、基金与专精特新企业代表；行政嘉宾按议题匹配协调，不做无法兑现的硬承诺。",
             13, 12)
    add_footer(s, 7)


def slide_calendar(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_title_bar(s, "07  2026下半年关键节点（建议版）",
                  "可按总经理日程与项目销售节奏微调；Excel表可直接改期")
    rows = [
        ["7–8月", "启动与暖场", "一滴水闭门沙龙（虹口）", "虹口项目", "锁档位、定片区轮动"],
        ["9月", "旗舰高光", "一滴水AI商业化落地公开日/晚宴", "虹口+黄浦露出", "最大品牌场"],
        ["10月", "产业走出去", "杨浦/虹口大厂参访周", "三区线索回流", "行政+交易双线"],
        ["11月", "片区轮动", "黄浦主题沙龙 + 青浦产业日", "黄浦/青浦项目", "服务去化节点"],
        ["12月", "收官复盘", "华东区三区联动私董宴", "区域总主持", "沉淀年度资产"],
    ]
    add_table(s, Inches(0.35), Inches(1.45), Inches(12.6), Inches(4.5),
              ["时间", "阶段", "建议动作", "项目联动", "对区域总动作"],
              rows, [1.2, 1.5, 3.8, 2.5, 2.6], body_size=12, header_size=12)
    add_footer(s, 8)


def slide_tech_path(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_title_bar(s, "08  大厂与 TMT 链接路径：从看见到交易",
                  "把杨浦、虹口的产业密度，变成绿城华东区的客群与合作入口")
    steps = [
        ("1 邀约", "峰会/沙龙定向邀约\n大厂业务线、创投、TMT CXO"),
        ("2 场内对话", "圆桌+案例路演\n绿城代表参与主持/致辞"),
        ("3 参访走访", "组织赴杨浦/虹口\n大厂与园区深度交流"),
        ("4 项目承接", "会后参观黄浦/\n虹口/青浦项目"),
        ("5 线索运营", "建立华东区线索池\n季度复盘与二次邀约"),
    ]
    for i, (t, b) in enumerate(steps):
        x = Inches(0.35) + i * Inches(2.55)
        add_rect(s, x, Inches(1.6), Inches(2.4), Inches(3.5), LIGHT_GREY)
        add_rect(s, x, Inches(1.6), Inches(2.4), Inches(0.7), GREEN)
        add_text(s, x + Inches(0.1), Inches(1.72), Inches(2.2), Inches(0.5),
                 t, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.15), Inches(2.55), Inches(2.1), Inches(2.3),
                 b, size=12, color=DARK, align=PP_ALIGN.CENTER)
    add_card(s, Inches(0.35), Inches(5.35), Inches(12.6), Inches(1.3),
             "可优先触达的方向（示例，按议题动态匹配）",
             "云计算与大模型应用、智能硬件、产业互联网、创投机构、专精特新与硬科技企业；目标不是“名单炫耀”，而是每次活动留下可跟进的企业/个人线索。",
             13, 12)
    add_footer(s, 9)


def slide_project_link(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_title_bar(s, "09  项目联动设计：服务去化，不抢项目总风头",
                  "区域总抓叙事与资源，项目总抓接待与转化")
    rows = [
        ["虹口项目（如潮鸣外滩等）", "一滴水主场后夜访/次日参观", "高端圈层晚宴+样板体验", "北外滩生活方式背书"],
        ["黄浦项目", "企业主沙龙后专场参观", "CBD客群包场、城市更新对话", "中心城区企业置业叙事"],
        ["青浦项目", "产业日/家庭日", "企业落位座谈+产品讲解", "新城置业与人才安居"],
    ]
    add_table(s, Inches(0.35), Inches(1.45), Inches(12.6), Inches(2.8),
              ["项目片区", "联动方式", "现场动作", "区域总可强调的点"],
              rows, [2.8, 3.2, 3.4, 3.2], body_size=12, header_size=12)
    add_card(s, Inches(0.35), Inches(4.5), Inches(6.1), Inches(2.1),
             "区域总角色建议",
             "开场定调片区战略、关键场次致辞、与行政/大厂嘉宾合影与闭门交流；不对单盘说过度承诺。",
             13, 12)
    add_card(s, Inches(6.65), Inches(4.5), Inches(6.3), Inches(2.1),
             "项目总角色建议",
             "接待动线、产品讲解、意向跟进；活动物料统一打“绿城华东区×XX项目”，强化片区母品牌。",
             13, 12)
    add_footer(s, 10)


def slide_tiers(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_title_bar(s, "10  赞助档位建议（金额可谈）",
                  "主推“华东区战略合作伙伴”打包，既有高度也有可执行颗粒度")
    rows = [
        ["钻石·华东区年度战略合作伙伴", "¥500,000", "下半年全系列主冠名", "区域总首选，片区共识最强"],
        ["铂金·片区首席合作伙伴", "¥300,000", "一滴水主场+2场片区联动", "性价比最高的主推档"],
        ["黄金·一滴水系列冠名", "¥150,000", "旗舰峰会/晚宴系列", "适合先做品牌高光"],
        ["银牌·项目专场赞助", "¥80,000", "单片区项目联动日", "项目总可分摊预算"],
        ["单场闭门沙龙/晚宴", "¥50,000–100,000", "单场冠名", "灵活试水"],
    ]
    add_table(s, Inches(0.35), Inches(1.4), Inches(12.6), Inches(4.2),
              ["档位", "建议金额", "覆盖范围", "沟通话术"],
              rows, [3.4, 1.8, 3.6, 3.8], body_size=12, header_size=12)
    add_text(s, Inches(0.45), Inches(5.85), Inches(12.4), Inches(0.8),
             "明日沟通建议：优先锁定铂金¥30万或钻石¥50万；若预算敏感，可先签黄金¥15万启动9月旗舰场，再追加片区轮动。",
             size=13, bold=True, color=DARK_GREEN)
    add_footer(s, 11)


def slide_rights(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_title_bar(s, "11  权益方案对照（核心权益摘要）",
                  "完整勾选矩阵见配套 Excel「权益矩阵」页")
    rows = [
        ["绿城华东区战略合作伙伴称号", "●", "●", "○", "—"],
        ["一滴水旗舰场冠名/联合冠名", "主冠名", "联合冠名", "主冠名", "—"],
        ["区域总致辞/对话席位", "●", "●", "○", "按场次"],
        ["大厂参访随行名额", "8–12人", "6–8人", "4人", "—"],
        ["三区项目联动露出", "全覆盖", "2区", "虹口", "1区"],
        ["现场展位/品牌陈列", "优先", "标准", "标准", "按场"],
        ["媒体与社群传播包", "全年包", "季度包", "单系列", "单场"],
        ["线索复盘会（对区域总）", "双月", "季度", "会后1次", "—"],
    ]
    add_table(s, Inches(0.3), Inches(1.35), Inches(12.7), Inches(5.2),
              ["权益项", "钻石50万", "铂金30万", "黄金15万", "银牌8万"],
              rows, [3.6, 2.2, 2.2, 2.2, 2.2], body_size=11, header_size=12)
    add_footer(s, 12)


def slide_package(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_title_bar(s, "12  主推打包：铂金·片区首席合作伙伴（¥30万）",
                  "最适合明日拍板：高度够、动作满、预算可解释")
    left = (
        "包含什么\n\n"
        "1）一滴水旗舰峰会/晚宴联合冠名 1场\n"
        "2）精品闭门沙龙 2场（虹口+黄浦或青浦）\n"
        "3）杨浦/虹口大厂或TMT参访 1次\n"
        "4）片区项目联动日 2场\n"
        "5）绿城华东区首席合作伙伴称号与传播包\n"
        "6）季度线索复盘会 2次（对区域总）"
    )
    right = (
        "对总经理的回报叙事\n\n"
        "· 客户：每场沉淀高净值/企业线索\n"
        "· 行政：有产业主题的对外接口\n"
        "· 片区：三区轮动，强化区域统筹\n"
        "· 品牌：AI×人居的差异化声量\n"
        "· 管理：Excel排期可下发项目总执行\n\n"
        "升级路径：执行满意后可补差升级钻石档"
    )
    add_card(s, Inches(0.4), Inches(1.45), Inches(6.2), Inches(5.15), "¥300,000 交付包", left, 15, 13)
    add_card(s, Inches(6.8), Inches(1.45), Inches(6.1), Inches(5.15), "沟通卖点", right, 15, 13)
    add_footer(s, 13)


def slide_outcomes(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_title_bar(s, "13  预期成果（下半年可衡量）",
                  "用区域总能听懂的指标说话")
    metrics = [
        ("6–8场", "高质量活动落地\n覆盖办行/办大众/参访/项目日"),
        ("200+", "有效圈层触达\n大厂/TMT/基金/企业主"),
        ("3区", "黄浦·虹口·青浦\n形成统一叙事与轮动"),
        ("双月复盘", "线索与行政互动\n沉淀为华东区资产"),
    ]
    for i, (n, d) in enumerate(metrics):
        x = Inches(0.4) + i * Inches(3.2)
        add_rect(s, x, Inches(1.6), Inches(3.05), Inches(3.3), LIGHT_GREY)
        add_text(s, x + Inches(0.15), Inches(2.0), Inches(2.75), Inches(0.8),
                 n, size=28, bold=True, color=DARK_GREEN, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.2), Inches(3.0), Inches(2.65), Inches(1.5),
                 d, size=13, color=DARK, align=PP_ALIGN.CENTER)
    add_card(s, Inches(0.4), Inches(5.15), Inches(12.5), Inches(1.45),
             "风险边界（体现专业度）",
             "嘉宾级别尽力匹配议题但不做“必到某人”硬承诺；行政出席按官方流程协调；个人敏感信息不做违规名单交付；项目转化由绿城销售体系闭环，我们负责场内质量与线索交接。",
             13, 12)
    add_footer(s, 14)


def slide_next(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_title_bar(s, "14  明日共识清单 & 推进节奏",
                  "带这三问进会议室，方便总经理当场拍板")
    qs = [
        ("问1", "下半年主推铂金30万，还是钻石50万？是否需要拆到项目分摊？"),
        ("问2", "首场定在一滴水（虹口）的哪个月？黄浦/青浦谁先轮动？"),
        ("问3", "区域总希望亲自站台的场次有几场？需要对接哪些行政条线？"),
    ]
    for i, (a, b) in enumerate(qs):
        y = Inches(1.45) + i * Inches(1.05)
        add_rect(s, Inches(0.45), y, Inches(12.4), Inches(0.9), LIGHT_GREY)
        add_rect(s, Inches(0.45), y, Inches(1.2), Inches(0.9), GOLD)
        add_text(s, Inches(0.55), y, Inches(1.0), Inches(0.9),
                 a, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.9), y + Inches(0.2), Inches(10.7), Inches(0.55),
                 b, size=14, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(1.7),
             "会后7日内：确认档位 → 出正式报价与排期 → 指定绿城对接人（区域办+项目）→ 启动首场筹备。\n\n"
             "协同方：上海市杨浦区科技企业联合会 / 见微知海新质商业生态\n"
             "对接人：胡继刚  13262607888",
             size=13, color=DARK)
    add_footer(s, 15)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide_cover(prs)
    slide_one_pager(prs)
    slide_insight(prs)
    slide_engines(prs)
    slide_map(prs)
    slide_venue(prs)
    slide_matrix(prs)
    slide_calendar(prs)
    slide_tech_path(prs)
    slide_project_link(prs)
    slide_tiers(prs)
    slide_rights(prs)
    slide_package(prs)
    slide_outcomes(prs)
    slide_next(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"saved: {OUT}")


if __name__ == "__main__":
    main()
