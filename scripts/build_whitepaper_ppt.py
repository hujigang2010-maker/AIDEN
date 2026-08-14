# -*- coding: utf-8 -*-
"""白皮书配套简报 PPT（研究会研究文稿第一号）。"""
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

OUT = "/workspace/whitepaper/面向2030新能源储能与新型电力系统白皮书-简报.pptx"
LOGO = "/workspace/whitepaper/assets/logo_fudan_hprc.png"
CHARTS = "/workspace/whitepaper/assets/charts"
BLUE = RGBColor(0x0E, 0x4E, 0x9B)
RED = RGBColor(0xC8, 0x10, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x5B, 0x61, 0x6B)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
TEAL = RGBColor(0x3A, 0xA1, 0x7E)
LIGHT = RGBColor(0xDC, 0xE6, 0xF1)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
PAGE = [0]


def add_bg(slide, color=WHITE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def bar(slide, y=0, h=0.08, color=BLUE):
    s = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(y), prs.slide_width, Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


def txt(slide, l, t, w, h, text, size=18, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, font="微软雅黑"):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return box


def bullets(slide, l, t, w, h, items, size=16, color=DARK, spacing=10):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing)
        p.level = 0
        run = p.add_run()
        run.text = "•  " + item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "微软雅黑"
    return box


def footer(slide):
    PAGE[0] += 1
    txt(slide, 0.5, 7.15, 10.2, 0.28,
        "复旦大学住房政策研究会 · 研究文稿第一号  FDU-HPRS-WP-2026-01",
        size=10, color=GRAY)
    txt(slide, 11.5, 7.15, 1.4, 0.28, str(PAGE[0]), size=10, color=GRAY,
        align=PP_ALIGN.RIGHT)


def img(slide, path, l, t, w):
    slide.shapes.add_picture(path, Inches(l), Inches(t), width=Inches(w))


def card(slide, l, t, w, h, fill=LIGHT):
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.adjustments[0] = 0.08
    return s


# 1 封面
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
bar(s, 0, 0.1)
img(s, LOGO, 0.65, 0.35, 6.6)
txt(s, 0.7, 2.35, 12, 0.4, "研究会研究文稿 · 第一号", 16, False, GRAY)
txt(s, 0.7, 2.9, 12, 0.7, "面向 2030", 22, True, GRAY)
txt(s, 0.7, 3.45, 12, 1.0, "新能源、储能与新型电力系统白皮书", 32, True, BLUE)
txt(s, 0.7, 4.55, 12, 0.5,
    "住房、社区与城市空间如何成为能源转型的关键节点", 18, False, GRAY)
txt(s, 0.7, 6.35, 12, 0.4,
    "FDU-HPRS-WP-2026-01    二〇二六年八月 · 上海", 14, False, GRAY)
bar(s, 7.42, 0.08, RED)
PAGE[0] += 1

# 2 核心命题
s = prs.slides.add_slide(BLANK)
bar(s)
txt(s, 0.6, 0.28, 12, 0.45, "核心命题：住房即能源节点", 26, True, BLUE)
txt(s, 0.6, 0.85, 12, 0.4,
    "新能源提供电量，储能提供时间，电力系统提供秩序，住房提供最后一公里的空间与同意。",
    15, False, GRAY)
bullets(s, 0.7, 1.45, 12, 5.3, [
    "2026年6月《新型能源体系建设“十五五”规划》、8月《新型电力系统建设“十五五”规划》把2030年钉死为“初步建成”验收年。",
    "约束已从“能不能装得上”转向“能不能用得稳、调得动、分得清、住得起”。",
    "屋顶产权、车位充电权、物业运营权、保障房绿色溢价、居民电费可负担性，已与电力系统重构深度咬合。",
    "本白皮书从住房政策视角给出十项判断、五类空间路径与十二条建议。",
], size=16, spacing=14)
footer(s)

# 3 十项判断（上）
s = prs.slides.add_slide(BLANK)
bar(s)
txt(s, 0.6, 0.25, 12, 0.4, "面向 2030 的十项核心判断（1—5）", 24, True, BLUE)
items = [
    ("01", "2030是初步建成，不是碳中和完成年"),
    ("02", "装机已绿，电量与调节结构尚未切换"),
    ("03", "项目竞争力取决于系统价值而非资源禀赋"),
    ("04", "储能从配建义务转向可交易调节资产"),
    ("05", "电网投资必须补上相对电源的欠账"),
]
for i, (n, t) in enumerate(items):
    y = 0.9 + i * 1.15
    card(s, 0.6, y, 12.1, 1.02)
    txt(s, 0.85, y + 0.22, 1.2, 0.55, n, 22, True, RED)
    txt(s, 2.2, y + 0.28, 10.2, 0.5, t, 18, True, DARK)
footer(s)

# 4 十项判断（下）
s = prs.slides.add_slide(BLANK)
bar(s)
txt(s, 0.6, 0.25, 12, 0.4, "面向 2030 的十项核心判断（6—10）", 24, True, BLUE)
items = [
    ("06", "用电侧协同升级为与电源、电网并列的第三支柱"),
    ("07", "住房是分布式能源最大未充分开发资产"),
    ("08", "电动汽车进入居住区，本质是移动储能入户"),
    ("09", "算力是增速最快、也最适合源网荷储一体的新负荷"),
    ("10", "社会可接受性最终在电费、住房品质与社区安全上兑现"),
]
for i, (n, t) in enumerate(items):
    y = 0.9 + i * 1.15
    card(s, 0.6, y, 12.1, 1.02)
    txt(s, 0.85, y + 0.22, 1.2, 0.55, n, 22, True, RED)
    txt(s, 2.2, y + 0.28, 10.2, 0.5, t, 18, True, DARK)
footer(s)

# 5 2030坐标
s = prs.slides.add_slide(BLANK)
bar(s)
txt(s, 0.5, 0.22, 12, 0.4, "2030 年验收清单：必须同时完成的目标", 24, True, BLUE)
img(s, os.path.join(CHARTS, "chart01_capacity_2025_2030.png"), 0.35, 0.75, 8.3)
bullets(s, 8.7, 1.1, 4.3, 5.5, [
    "风光装机 28 亿千瓦以上",
    "非化石电量占比 50%",
    "新型储能 3 亿千瓦",
    "抽水蓄能约 1.6 亿千瓦",
    "西电东送超过 4.2 亿千瓦",
    "配电网支撑 9 亿千瓦分布式",
    "虚拟电厂调节 >5000 万千瓦",
    "充电网络支撑 1.1 亿辆以上电动汽车",
    "全国统一电力市场基本建成",
], size=14, spacing=8)
footer(s)

# 6 新能源
s = prs.slides.add_slide(BLANK)
bar(s)
txt(s, 0.5, 0.22, 12, 0.4, "新能源：电量主体如何成为可调节主力", 24, True, BLUE)
img(s, os.path.join(CHARTS, "chart03_wind_solar_mix.png"), 0.3, 0.8, 7.6)
bullets(s, 8.1, 1.0, 4.8, 5.8, [
    "2025 年风光装机 18.4 亿千瓦，历史性超过火电。",
    "可再生能源电量约 38%，距 50% 仍差约 12 个百分点。",
    "分布式光伏 5.3 亿千瓦，已近光伏一半。",
    "利用率分档：地区 85%—95%，全国约 90%。",
    "主题词从“规模替代”转为“可靠替代”。",
], size=15, spacing=12)
footer(s)

# 7 储能
s = prs.slides.add_slide(BLANK)
bar(s)
txt(s, 0.5, 0.22, 12, 0.4, "储能：从强制配建到系统压舱石", 24, True, BLUE)
img(s, os.path.join(CHARTS, "chart04_storage_growth.png"), 0.25, 0.75, 6.6)
img(s, os.path.join(CHARTS, "chart05_storage_duration.png"), 6.85, 0.85, 6.2)
footer(s)

# 8 电力系统
s = prs.slides.add_slide(BLANK)
bar(s)
txt(s, 0.5, 0.22, 12, 0.4, "电力系统：主配微协同，配网才是住房侧主战场", 22, True, BLUE)
img(s, os.path.join(CHARTS, "chart06_west_to_east.png"), 0.3, 0.8, 6.4)
bullets(s, 6.9, 0.95, 5.9, 5.8, [
    "主网：西电东送超 4.2 亿千瓦，东部不能把保供外包给西部。",
    "配网：支撑 9 亿千瓦分布式接入；老旧小区、单电源高层、城中村是底线工程。",
    "微网：医院、保障房、防灾中心“能孤岛”。",
    "市场：给小而散的住房侧资源留入口，避免统一市场变成大用户市场。",
], size=15, spacing=12)
footer(s)

# 9 住房节点
s = prs.slides.add_slide(BLANK)
bar(s)
txt(s, 0.5, 0.22, 12, 0.4, "住房即能源节点：四层结构", 24, True, BLUE)
img(s, os.path.join(CHARTS, "chart07_housing_energy_node.png"), 0.7, 0.75, 11.8)
footer(s)

# 10 五类住房
s = prs.slides.add_slide(BLANK)
bar(s)
txt(s, 0.5, 0.22, 12, 0.4, "五类住房空间，五条差异化路径", 24, True, BLUE)
img(s, os.path.join(CHARTS, "chart08_housing_types.png"), 0.25, 0.7, 7.5)
bullets(s, 7.9, 0.95, 4.9, 5.8, [
    "新建商品房：能源就绪交付。",
    "存量小区：公共屋顶收益进业主账户。",
    "老旧/城中村/单电源：先配电与安全，再光伏。",
    "保障房：政府业主，最适合制度试验。",
    "农村：从屋顶租赁升级为光储直柔农宅。",
], size=15, spacing=12)
footer(s)

# 11 新业态
s = prs.slides.add_slide(BLANK)
bar(s)
txt(s, 0.5, 0.22, 12, 0.4, "新业态：虚拟电厂 · 车网互动 · 绿电直连 · 算电协同", 22, True, BLUE)
img(s, os.path.join(CHARTS, "chart09_new_business.png"), 0.25, 0.7, 7.4)
bullets(s, 7.85, 0.95, 5.0, 5.8, [
    "虚拟电厂投资约为传统电厂 10%—20%。",
    "居住区慢充才是可调度容量的主力。",
    "有序充电默认开通，V2G 自愿加入。",
    "绿电直连须支付系统备用，防止专线特权。",
    "算力与电力、住房配套同单元规划。",
], size=15, spacing=12)
footer(s)

# 12 路线图
s = prs.slides.add_slide(BLANK)
bar(s)
txt(s, 0.5, 0.22, 12, 0.4, "三条时间线：2026—2030", 24, True, BLUE)
img(s, os.path.join(CHARTS, "chart10_roadmap.png"), 0.55, 0.7, 12.2)
footer(s)

# 13 十二条建议（上）
s = prs.slides.add_slide(BLANK)
bar(s)
txt(s, 0.5, 0.22, 12, 0.4, "十二条政策建议（1—6）", 24, True, BLUE)
left = [
    "1. 系统价值作为新能源核准首要判据",
    "2. 容量电价不向低收入居民无差别传导",
    "3. 按调节时长付费，打开长时储能闭环",
    "4. 配电网更新列入城市更新基础清单",
    "5. 立法明确屋顶等公共部位能源使用权",
    "6. 保障房率先成为光储直柔与虚拟电厂样板",
]
bullets(s, 0.8, 1.0, 11.5, 5.8, left, size=18, spacing=16)
footer(s)

# 14 十二条建议（下）
s = prs.slides.add_slide(BLANK)
bar(s)
txt(s, 0.5, 0.22, 12, 0.4, "十二条政策建议（7—12）", 24, True, BLUE)
right = [
    "7. 居住区充电：有序充电默认开通、V2G 自愿加入",
    "8. 虚拟电厂与住房侧聚合资源获得对等市场入口",
    "9. 绿电直连必须支付系统备用",
    "10. 算力设施与电力、住房配套同单元规划",
    "11. 社区储能安全准入与保险强制挂钩",
    "12. 建立住房能源账单与极端天气对冲机制",
]
bullets(s, 0.8, 1.0, 11.5, 5.8, right, size=18, spacing=16)
footer(s)

# 15 结语
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
bar(s, 0, 0.1)
txt(s, 0.7, 1.6, 12, 0.5, "结语", 18, False, GRAY)
txt(s, 0.7, 2.2, 12, 1.4,
    "2030 年的验收，在变压器，也在电梯。", 28, True, BLUE)
txt(s, 0.7, 3.7, 11.5, 1.6,
    "远望 2030，不是远望一个更绿的装机表，而是远望一套更公平、更安全、可负担的居住与用能秩序。",
    18, False, DARK)
txt(s, 0.7, 5.7, 12, 0.4,
    "复旦大学住房政策研究会  ·  FDU-HPRS-WP-2026-01", 14, False, GRAY)
bar(s, 7.42, 0.08, RED)
PAGE[0] += 1

prs.save(OUT)
print("saved", OUT, "slides", len(prs.slides))
