#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""生成《许俊贤技能清单》16:9 可编辑 PPTX。

来源：许俊贤聊天记录里甩出的 13 条链接（12 个产品 + 1 个在线入口）。
口径：按产物分赛道，不按 Star 排名。Star 抓取于 2026-08-23。
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

INK = RGBColor(0x11, 0x11, 0x11)
PAPER = RGBColor(0xF4, 0xF0, 0xE6)
PAPER2 = RGBColor(0xE8, 0xE2, 0xD4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xE3, 0x1C, 0x24)
GRAY = RGBColor(0x6B, 0x6B, 0x6B)
MUTED = RGBColor(0x4A, 0x4A, 0x4A)
LINE = RGBColor(0xD6, 0xD0, 0xC2)
CREAM = RGBColor(0xFB, 0xF8, 0xF1)
FOOT_GRAY = RGBColor(0x9A, 0x96, 0x8C)

FONT = "Microsoft YaHei"
PAGE_W = Inches(13.333)
PAGE_H = Inches(7.5)
FOOTER = "许俊贤技能清单  ·  按产物分赛道  ·  2026.08.23"
TOTAL = 19

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "output"
OUT = OUT_DIR / "许俊贤技能清单.pptx"


def _set_run_font(run, size, color, bold=False, italic=False, name=FONT):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def add_text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.12):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    first = True
    for para in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = line_spacing
        runs = para if (para and isinstance(para[0], tuple)) else [para]
        for item in runs:
            text, size, color, bold = item[0], item[1], item[2], item[3]
            italic = item[4] if len(item) > 4 else False
            r = p.add_run()
            r.text = text
            _set_run_font(r, size, color, bold, italic)
    return tb


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75,
             shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def footer(slide, idx, dark=False):
    color = RGBColor(0x88, 0x88, 0x88) if dark else FOOT_GRAY
    add_text(slide, Inches(0.45), Inches(7.18), Inches(10.2), Inches(0.22),
             [[(FOOTER, 9, color, False)]])
    add_text(slide, Inches(11.4), Inches(7.18), Inches(1.45), Inches(0.22),
             [[(f"{idx:02d}  /  {TOTAL:02d}", 9, color, False)]],
             align=PP_ALIGN.RIGHT)


def kicker_bar(slide, kicker, title, subtitle=None):
    set_bg(slide, PAPER)
    add_rect(slide, 0, 0, PAGE_W, Inches(0.08), RED)
    add_text(slide, Inches(0.5), Inches(0.22), Inches(12.2), Inches(0.28),
             [[(kicker, 11, RED, True)]])
    add_text(slide, Inches(0.5), Inches(0.48), Inches(12.2), Inches(0.5),
             [[(title, 26, INK, True)]])
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.00), Inches(12.2), Inches(0.32),
                 [[(subtitle, 12, GRAY, False)]])


def hairline(slide, x, y, w, color=LINE):
    add_rect(slide, x, y, w, Pt(1.0), color)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---------------------------------------------------------------------------
# 数据
# ---------------------------------------------------------------------------

DECK_SKILLS = [
    {
        "id": "01",
        "name": "ppt-master",
        "who": "花叔Hugo",
        "star": "48,747",
        "out": "可编辑 PPTX",
        "one": "真正的原生 PowerPoint 对象模型",
        "fit": "领导改稿、客户交付",
        "cmd": "npx skills add hugohe3/ppt-master",
        "url": "github.com/hugohe3/ppt-master",
    },
    {
        "id": "02",
        "name": "frontend-slides",
        "who": "张咋啦",
        "star": "27,997",
        "out": "单文件 HTML",
        "one": "视觉上限最高的网页演示",
        "fit": "技术分享、在线传播",
        "cmd": "npx skills add zarazhangrui/frontend-slides",
        "url": "github.com/zarazhangrui/frontend-slides",
    },
    {
        "id": "03",
        "name": "guizang-ppt-skill",
        "who": "归藏",
        "star": "24,671",
        "out": "单文件 HTML",
        "one": "瑞士风 / 杂志风，自带快捷键",
        "fit": "线下分享、观点表达",
        "cmd": "npx skills add op7418/guizang-ppt-skill",
        "url": "github.com/op7418/guizang-ppt-skill",
    },
    {
        "id": "04",
        "name": "baoyu-slide-deck",
        "who": "宝玉",
        "star": "25,263*",
        "out": "图片卡片",
        "one": "为阅读转发设计，不是为投影",
        "fit": "课件、社媒整页图",
        "cmd": "npx skills add JimLiu/baoyu-skills --skill baoyu-slide-deck",
        "url": "github.com/JimLiu/baoyu-skills",
    },
    {
        "id": "05",
        "name": "huashu-design",
        "who": "花叔",
        "star": "23,423",
        "out": "HTML + PPTX",
        "one": "工作室流程：先三方向再展开",
        "fit": "品牌发布、要设计流程",
        "cmd": "npx skills add alchaincyf/huashu-design",
        "url": "github.com/alchaincyf/huashu-design",
    },
    {
        "id": "06",
        "name": "html-ppt-skill",
        "who": "Lewis",
        "star": "8,024",
        "out": "单文件 HTML",
        "one": "计时器 + 逐字稿，控场优先",
        "fit": "上台要讲稿和计时",
        "cmd": "npx skills add lewislulu/html-ppt-skill",
        "url": "github.com/lewislulu/html-ppt-skill",
    },
    {
        "id": "07",
        "name": "qiaomu-anything-to-notebooklm",
        "who": "乔木",
        "star": "5,783",
        "out": "多格式卡片",
        "one": "先把素材变成结构，不是终稿",
        "fit": "公众号/PDF/播客初整理",
        "cmd": "npx skills add joeseesun/qiaomu-anything-to-notebooklm",
        "url": "github.com/joeseesun/qiaomu-anything-to-notebooklm",
    },
]

ALL_ITEMS = [
    ("演示", "ppt-master", "花叔Hugo", "可编辑 PPTX", "48,747", "MIT"),
    ("演示", "frontend-slides", "张咋啦", "单文件 HTML", "27,997", "MIT"),
    ("演示", "baoyu-slide-deck", "宝玉", "图片卡片", "25,263*", "MIT"),
    ("演示", "guizang-ppt-skill", "归藏", "单文件 HTML", "24,671", "AGPL-3.0"),
    ("演示", "huashu-design", "花叔", "HTML + PPTX", "23,423", "MIT"),
    ("演示", "html-ppt-skill", "Lewis", "单文件 HTML", "8,024", "MIT"),
    ("演示", "qiaomu-anything-to-notebooklm", "乔木", "多格式卡片", "5,783", "MIT"),
    ("模板", "beautiful-html-templates", "张咋啦", "34 套 HTML 模板", "4,324", "MIT"),
    ("白板", "beautiful-feishu-whiteboard", "张咋啦", "35 套飞书色板", "681", "MIT"),
    ("白板", "larkboard-graphy", "F俊W", "叙事图表白板", "14", "MIT"),
    ("改稿", "htmledit", "财猫", "可视化改 HTML", "8", "CC-BY-NC"),
    ("科研", "XRD-SKILL", "LucianaiB2004", "XRD 一图流白板", "3", "MIT"),
]


# ---------------------------------------------------------------------------
# 页面
# ---------------------------------------------------------------------------

def slide_cover(prs):
    s = blank(prs)
    set_bg(s, INK)
    add_rect(s, 0, 0, Inches(0.18), PAGE_H, RED)
    add_text(s, Inches(0.7), Inches(0.85), Inches(11.5), Inches(0.35),
             [[("许俊贤聊天记录整理  ·  不是排行榜", 13, RED, True)]])
    add_text(s, Inches(0.7), Inches(1.45), Inches(12), Inches(1.7),
             [[("一份被甩过来的", 36, WHITE, True)],
              [("视觉生产技能清单", 36, WHITE, True)]],
             line_spacing=1.05)
    add_rect(s, Inches(0.7), Inches(3.45), Inches(1.6), Inches(0.06), RED)
    add_text(s, Inches(0.7), Inches(3.75), Inches(11.5), Inches(1.3),
             [[("十二件工具，四条赛道：演示稿、HTML 模板、改稿器、飞书白板。", 16, PAPER, False)],
              [("上一份材料评的是「七款中文 AI PPT」。这份把地图扩到生成之后怎么改、怎么上白板。", 16, PAPER, False)]],
             line_spacing=1.28)
    stats = [("12", "件产品"), ("13", "条链接"), ("4", "条赛道"), ("2026.08.23", "Star 抓取")]
    for i, (n, k) in enumerate(stats):
        x = Inches(0.7) + i * Inches(3.05)
        add_text(s, x, Inches(5.45), Inches(2.8), Inches(0.45),
                 [[(n, 22, WHITE, True)]])
        add_text(s, x, Inches(5.95), Inches(2.8), Inches(0.3),
                 [[(k, 12, FOOT_GRAY, False)]])
    add_notes(s, "开场：这不是锐评排名，是许俊贤甩过来的技能地图。四条赛道比十二个名字更重要。")


def slide_source(prs, idx):
    s = blank(prs)
    kicker_bar(s, "01  来源", "聊天记录里的五段话",
               "按原话分组还原。htmledit 仓库和 htmledit.ai 算同一件产品的两个入口。")
    bundles = [
        ("01", "张咋啦三件套 + 宝玉",
         "baoyu-slide-deck  ·  frontend-slides  ·  beautiful-feishu-whiteboard  ·  beautiful-html-templates",
         "演示 + 模板 + 白板，一个人覆盖三条赛道。"),
        ("02", "归藏 / 乔木 / 花叔 / PPT Master",
         "guizang-ppt-skill  ·  qiaomu-anything-to-notebooklm  ·  huashu-design  ·  ppt-master",
         "七款 PPT 横评的主体，都在这一段。"),
        ("03", "Lewis + 财猫改稿",
         "html-ppt-skill  ·  hamutama/htmledit  ·  htmledit.ai",
         "HTML 演示讲完之后，补上「不会写代码也能改网页」。"),
        ("04", "科研一图流",
         "LucianaiB2004/XRD-SKILL",
         "主打用画板做科研图表，不是通用 PPT。"),
        ("05", "数据叙事白板  ·  F俊W",
         "lostvita/larkboard-graphy",
         "三段式架构 + Widgets 叙事层，把普通图表变成带洞察的汇报图。"),
    ]
    for i, (num, title, urls, note) in enumerate(bundles):
        y = Inches(1.42) + i * Inches(1.05)
        add_rect(s, Inches(0.5), y, Inches(12.35), Inches(0.95), WHITE)
        add_text(s, Inches(0.7), y + Inches(0.08), Inches(0.6), Inches(0.35),
                 [[(num, 16, RED, True)]])
        add_text(s, Inches(1.4), y + Inches(0.08), Inches(11.1), Inches(0.32),
                 [[(title, 15, INK, True)]])
        add_text(s, Inches(1.4), y + Inches(0.40), Inches(11.1), Inches(0.22),
                 [[(urls, 11, MUTED, False)]])
        add_text(s, Inches(1.4), y + Inches(0.62), Inches(11.1), Inches(0.25),
                 [[(note, 12, GRAY, False)]])
    footer(s, idx)
    add_notes(s, "还原聊天结构，方便对照原话。财猫的仓库和网站算一件产品。")


def slide_how(prs, idx):
    s = blank(prs)
    kicker_bar(s, "02  口径", "先问产物，再问场景",
               "Star 是人气，不是质量。宝玉的数字是整个 baoyu-skills 仓库。")
    cards = [
        ("产物", "交出去的文件是什么",
         "PPTX、单文件 HTML、整页图片、飞书可编辑白板、还是打开浏览器就能改的网页。"),
        ("场景", "第二天要交给谁",
         "领导改稿、自己上台、发群转发、飞书里继续画、实验室交图——人不同，工具不同。"),
        ("位置", "它在生产链的哪一环",
         "生成演示、选视觉模板、改已有 HTML、把数据画上白板。四环不要拿来互相比。"),
        ("协议", "能不能商用、怎么装",
         "多数 MIT；归藏 AGPL-3.0；财猫 CC-BY-NC-4.0。安装命令见第 15 页。"),
    ]
    for i, (k, t, b) in enumerate(cards):
        col, row = i % 2, i // 2
        x = Inches(0.5) + col * Inches(6.35)
        y = Inches(1.5) + row * Inches(2.55)
        add_rect(s, x, y, Inches(6.15), Inches(2.35), WHITE)
        add_rect(s, x, y, Inches(0.1), Inches(2.35), RED)
        add_text(s, x + Inches(0.35), y + Inches(0.25), Inches(5.5), Inches(0.3),
                 [[(k, 12, RED, True)]])
        add_text(s, x + Inches(0.35), y + Inches(0.6), Inches(5.5), Inches(0.45),
                 [[(t, 18, INK, True)]])
        add_text(s, x + Inches(0.35), y + Inches(1.2), Inches(5.5), Inches(0.85),
                 [[(b, 14, MUTED, False)]], line_spacing=1.25)
    footer(s, idx)
    add_notes(s, "读这份清单的方法：产物 → 场景 → 生产链位置。不要用 Star 给白板和 PPT 排名。")


def slide_routes(prs, idx):
    s = blank(prs)
    kicker_bar(s, "03  地图", "四条赛道，不要混着比")
    routes = [
        ("A  演示稿", "7 件",
         "把主题做成能讲、能发、能改的幻灯片。",
         "PPTX / HTML / 图片",
         "ppt-master · frontend-slides · 归藏 · 宝玉 · 花叔 · Lewis · 乔木"),
        ("B  HTML 模板", "1 件",
         "不从零生成，先选一套视觉系统再填内容。",
         "34 套可克隆模板",
         "beautiful-html-templates（张咋啦）"),
        ("C  改稿器", "1 件",
         "AI 生成之后，像改文档一样改文字、图、结构。",
         "浏览器本地编辑",
         "htmledit / htmledit.ai（财猫）"),
        ("D  飞书白板", "3 件",
         "可编辑看板、叙事图表、科研 XRD 一图流。",
         "飞书原生对象",
         "张咋啦白板 · larkboard-graphy · XRD-SKILL"),
    ]
    for i, (title, n, desc, out, who) in enumerate(routes):
        y = Inches(1.42) + i * Inches(1.32)
        add_rect(s, Inches(0.5), y, Inches(12.35), Inches(1.2), WHITE)
        add_rect(s, Inches(0.5), y, Inches(0.1), Inches(1.2), RED)
        add_text(s, Inches(0.85), y + Inches(0.12), Inches(7.5), Inches(0.35),
                 [[(title, 18, INK, True)]])
        add_text(s, Inches(8.6), y + Inches(0.14), Inches(3.9), Inches(0.32),
                 [[(n + "   ·   " + out, 12, RED, True)]], align=PP_ALIGN.RIGHT)
        add_text(s, Inches(0.85), y + Inches(0.5), Inches(11.6), Inches(0.28),
                 [[(desc, 13, MUTED, False)]])
        add_text(s, Inches(0.85), y + Inches(0.8), Inches(11.6), Inches(0.28),
                 [[(who, 12, GRAY, False)]])
    footer(s, idx)
    add_notes(s, "四条赛道是整份材料的骨架。后面每一页都挂在其中一条上。")


def slide_inventory(prs, idx):
    s = blank(prs)
    kicker_bar(s, "04  全表", "十二件产品，按赛道排列",
               "* 宝玉 Star 是整个 baoyu-skills 仓库。财猫另有在线入口 htmledit.ai。")
    headers = ["赛道", "名称", "作者", "产物", "Star", "协议"]
    col_w = [Inches(1.15), Inches(3.85), Inches(1.9), Inches(2.35),
             Inches(1.35), Inches(1.55)]
    x0, y0, rh = Inches(0.5), Inches(1.42), Inches(0.42)
    add_rect(s, x0, y0, Inches(12.15), Inches(0.38), INK)
    x = x0
    for i, h in enumerate(headers):
        add_text(s, x + Inches(0.08), y0, col_w[i] - Inches(0.1), Inches(0.38),
                 [[(h, 11, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[i]
    for r, row in enumerate(ALL_ITEMS):
        y = y0 + Inches(0.38) + r * rh
        bg = WHITE if r % 2 == 0 else PAPER2
        add_rect(s, x0, y, Inches(12.15), rh, bg)
        x = x0
        for c, val in enumerate(row):
            color = RED if c == 0 else (INK if c == 1 else MUTED)
            add_text(s, x + Inches(0.08), y, col_w[c] - Inches(0.1), rh,
                     [[(val, 11, color, c in (0, 1))]],
                     anchor=MSO_ANCHOR.MIDDLE)
            x += col_w[c]
    footer(s, idx)
    add_notes(s, "全表按赛道而不是按 Star。演示赛道七件占了人气，但白板和改稿解决的是另一类问题。")


def slide_deck_grid(prs, idx):
    s = blank(prs)
    kicker_bar(s, "05  赛道 A", "演示七款：一句话认人",
               "细节仍以「七款中文 AI PPT Skill 横评」为准。这里只作地图上的锚点。")
    for i, sk in enumerate(DECK_SKILLS):
        col, row = i % 4, i // 4
        # 7 items: 4 + 3, last row centered-ish by starting at 0.5
        if row == 1:
            x = Inches(0.5) + col * Inches(4.15)
            w = Inches(4.0)
        else:
            x = Inches(0.5) + col * Inches(3.15)
            w = Inches(3.0)
        y = Inches(1.45) + row * Inches(2.7)
        h = Inches(2.5)
        add_rect(s, x, y, w, h, WHITE)
        add_rect(s, x, y, w, Inches(0.08), RED)
        add_text(s, x + Inches(0.15), y + Inches(0.18), w - Inches(0.3), Inches(0.28),
                 [[(sk["id"] + "  " + sk["who"], 11, RED, True)]])
        add_text(s, x + Inches(0.15), y + Inches(0.48), w - Inches(0.3), Inches(0.55),
                 [[(sk["name"], 13, INK, True)]], line_spacing=1.05)
        add_text(s, x + Inches(0.15), y + Inches(1.05), w - Inches(0.3), Inches(0.55),
                 [[(sk["one"], 12, MUTED, False)]], line_spacing=1.15)
        add_text(s, x + Inches(0.15), y + Inches(1.7), w - Inches(0.3), Inches(0.28),
                 [[(sk["out"] + "  ·  " + sk["star"], 11, GRAY, False)]])
        add_text(s, x + Inches(0.15), y + Inches(2.05), w - Inches(0.3), Inches(0.28),
                 [[(sk["fit"], 12, INK, True)]])
    footer(s, idx)
    add_notes(s, "七款演示 Skill 压缩成两行卡片。需要细读请看上一份 18 页横评。")


def slide_deck_pick(prs, idx):
    s = blank(prs)
    kicker_bar(s, "06  赛道 A", "演示稿只问一句话：明天交给谁")
    rows = [
        ("投资人 / 领导要改", "ppt-master", "原生 PPTX，形状图表母版都能继续改。"),
        ("品牌发布、设计很挑", "huashu-design", "先出三个方向，再要 HTML 和可编辑 PPTX。"),
        ("技术分享、视觉上限", "frontend-slides", "三张预览定风格，单文件 HTML 发出去。"),
        ("线下瑞士风 / 杂志风", "guizang-ppt-skill", "快捷键翻页，双击浏览器开讲。"),
        ("上台要计时和讲稿", "html-ppt-skill", "演讲者窗口：当前页、下一页、逐字稿、计时器。"),
        ("课件 / 社媒整页卡片", "baoyu-slide-deck", "为阅读和转发设计，不是为投影。"),
        ("先把素材变成结构", "qiaomu-anything-to-notebooklm", "初稿和多格式转化，不是终稿。"),
    ]
    for i, (scene, name, why) in enumerate(rows):
        y = Inches(1.42) + i * Inches(0.74)
        add_rect(s, Inches(0.5), y, Inches(12.35), Inches(0.66), WHITE)
        add_rect(s, Inches(0.5), y, Inches(0.1), Inches(0.66), RED)
        add_text(s, Inches(0.8), y, Inches(4.2), Inches(0.66),
                 [[(scene, 14, INK, True)]], anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(5.15), y, Inches(3.4), Inches(0.66),
                 [[(name, 13, RED, True)]], anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(8.6), y, Inches(4.05), Inches(0.66),
                 [[(why, 12, MUTED, False)]], anchor=MSO_ANCHOR.MIDDLE)
    footer(s, idx)
    add_notes(s, "演示赛道选型：交给谁，就决定用哪一款。")


def slide_zhang_family(prs, idx):
    s = blank(prs)
    kicker_bar(s, "07  作者地图", "张咋啦一个人，三条赛道",
               "聊天记录第一段的核心。三件套不要互相替代，要串联。")
    trio = [
        ("frontend-slides", "演示稿", "27,997 Star",
         "从主题生成完整网页演示。12 套安全预设 + 34 套大胆模板，先出三张预览再开工。",
         "适合：技术分享、公开演讲。不适合：必须丢进 PowerPoint。"),
        ("beautiful-html-templates", "模板库", "4,324 Star",
         "34 套可克隆 HTML 模板。Agent 先问场合和气质，挑 3 个封面预览，等人选定再填全文。",
         "适合：已有内容、要挑视觉系统。不适合：从零长论证。"),
        ("beautiful-feishu-whiteboard", "飞书白板", "681 Star",
         "35 套经过实机验证的色板。产出是飞书文档里可继续改的原生白板，不是截图。",
         "适合：架构图、流程看板。不适合：当数据叙事图表引擎。"),
    ]
    for i, (name, lane, star, body, fit) in enumerate(trio):
        x = Inches(0.5) + i * Inches(4.2)
        add_rect(s, x, Inches(1.45), Inches(4.0), Inches(5.4), WHITE)
        add_rect(s, x, Inches(1.45), Inches(4.0), Inches(0.08), RED)
        add_text(s, x + Inches(0.22), Inches(1.7), Inches(3.55), Inches(0.28),
                 [[(lane, 12, RED, True)]])
        add_text(s, x + Inches(0.22), Inches(2.05), Inches(3.55), Inches(0.7),
                 [[(name, 16, INK, True)]], line_spacing=1.05)
        add_text(s, x + Inches(0.22), Inches(2.75), Inches(3.55), Inches(0.3),
                 [[(star, 12, GRAY, False)]])
        add_text(s, x + Inches(0.22), Inches(3.2), Inches(3.55), Inches(1.7),
                 [[(body, 13, MUTED, False)]], line_spacing=1.22)
        hairline(s, x + Inches(0.22), Inches(5.05), Inches(3.55))
        add_text(s, x + Inches(0.22), Inches(5.2), Inches(3.55), Inches(1.35),
                 [[(fit, 12, INK, False)]], line_spacing=1.22)
    footer(s, idx)
    add_notes(s, "张咋啦三件套：生成演示、选模板、上白板。frontend-slides 内部也会用到那 34 套模板。")


def slide_html_templates(prs, idx):
    s = blank(prs)
    kicker_bar(s, "08  赛道 B", "beautiful-html-templates",
               "不是生成器，是给 Agent 用的视觉系统图书馆。")
    points = [
        ("先问再挑", "强制问场合和气质，再从 index.json 里匹配 mood / tone / best_for / formality。"),
        ("三封面预览", "先做三套封面 HTML，等人用眼睛选，而不是用形容词猜。"),
        ("克隆再填", "选定后整夹克隆，按 preserve / replace / extend 规则改内容，不动视觉系统。"),
        ("和 frontend-slides 的关系", "frontend-slides 是完整演示工作流；本库是可复用模板底盘，也可被其他 Agent 单独调用。"),
    ]
    for i, (t, b) in enumerate(points):
        y = Inches(1.45) + i * Inches(1.2)
        add_rect(s, Inches(0.5), y, Inches(12.35), Inches(1.08), WHITE)
        add_rect(s, Inches(0.5), y, Inches(0.1), Inches(1.08), RED)
        add_text(s, Inches(0.85), y + Inches(0.14), Inches(11.7), Inches(0.32),
                 [[(t, 16, INK, True)]])
        add_text(s, Inches(0.85), y + Inches(0.5), Inches(11.7), Inches(0.45),
                 [[(b, 13, MUTED, False)]])
    footer(s, idx)
    add_notes(s, "模板库的价值是可预测的审美，不是一次生成。安装方式是克隆仓库并让 Agent 读 AGENTS.md。")


def slide_htmledit(prs, idx):
    s = blank(prs)
    kicker_bar(s, "09  赛道 C", "财猫 htmledit：生成之后怎么改",
               "仓库 github.com/hamutama/htmledit  ·  在线 htmledit.ai  ·  协议 CC-BY-NC-4.0")
    left = [
        ("问题", "AI 已经很容易生成网页和 HTML 演示。真正的门槛是改几个字、换一张图、挪一个段落。"),
        ("做法", "导入单个 HTML 或项目文件夹，点选即改文字、图片、链接、列表、表格。交互和翻页仍可用。"),
        ("隐私", "local-first：文件在浏览器本地读，草稿存在本机。清浏览器数据会丢草稿，重要修改要导出。"),
        ("入口", "打开 htmledit.ai 即可，不用注册。也可本地静态服务器，或装 Chrome 扩展（准备上架）。"),
    ]
    for i, (k, b) in enumerate(left):
        y = Inches(1.45) + i * Inches(1.2)
        add_rect(s, Inches(0.5), y, Inches(8.15), Inches(1.08), WHITE)
        add_text(s, Inches(0.7), y + Inches(0.12), Inches(7.75), Inches(0.28),
                 [[(k, 13, RED, True)]])
        add_text(s, Inches(0.7), y + Inches(0.42), Inches(7.75), Inches(0.55),
                 [[(b, 13, MUTED, False)]], line_spacing=1.18)
    add_rect(s, Inches(8.85), Inches(1.45), Inches(4.0), Inches(4.85), INK)
    add_text(s, Inches(9.1), Inches(1.7), Inches(3.55), Inches(0.35),
             [[("适合", 12, RED, True)]])
    add_text(s, Inches(9.1), Inches(2.1), Inches(3.55), Inches(1.5),
             [[("改 frontend-slides / 归藏 / Lewis / 花叔吐出的 HTML。不会写代码的同事也能动。", 13, WHITE, False)]],
             line_spacing=1.22)
    add_text(s, Inches(9.1), Inches(3.7), Inches(3.55), Inches(0.35),
             [[("不适合", 12, RED, True)]])
    add_text(s, Inches(9.1), Inches(4.1), Inches(3.55), Inches(1.8),
             [[("当 Skill 从零生成演示。复杂前端框架页可能改不完整。非商业协议，商用需另谈。", 13, PAPER, False)]],
             line_spacing=1.22)
    footer(s, idx)
    add_notes(s, "htmledit 补的是生成之后的人机改稿。它不是 PPT Skill，是 HTML 演示路线的下游。")


def slide_board_overview(prs, idx):
    s = blank(prs)
    kicker_bar(s, "10  赛道 D", "飞书白板三件，解决三种图",
               "共同前提：飞书账号 + 已登录的 lark-cli / whiteboard-cli。产出都是可继续改的原生对象。")
    cards = [
        ("审美看板", "beautiful-feishu-whiteboard",
         "35 套色板，从克制到大胆。Agent 构图，模板只管颜色和气质。",
         "架构图、流程、概念看板",
         "不是自动图表工具。"),
        ("叙事图表", "larkboard-graphy",
         "先分析数据再提 Widget 策略。标题区结论先行，图表区 + 信息区讲完一个故事。",
         "汇报图、经营分析、带洞察的柱线环",
         "不是通用看板美化。"),
        ("科研一图流", "XRD-SKILL",
         "二维 XRD 原始数据 → SVG/PNG/OpenAPI → 飞书可编辑画板。脚本主导，可重复跑。",
         "XRD 汇报图、实验对比、画板修复",
         "不是通用科研绘图套件。"),
    ]
    for i, (tag, name, body, fit, notfit) in enumerate(cards):
        x = Inches(0.5) + i * Inches(4.2)
        add_rect(s, x, Inches(1.45), Inches(4.0), Inches(5.4), WHITE)
        add_rect(s, x, Inches(1.45), Inches(4.0), Inches(0.08), RED)
        add_text(s, x + Inches(0.22), Inches(1.7), Inches(3.55), Inches(0.28),
                 [[(tag, 12, RED, True)]])
        add_text(s, x + Inches(0.22), Inches(2.05), Inches(3.55), Inches(0.85),
                 [[(name, 16, INK, True)]], line_spacing=1.05)
        add_text(s, x + Inches(0.22), Inches(2.95), Inches(3.55), Inches(1.55),
                 [[(body, 13, MUTED, False)]], line_spacing=1.22)
        hairline(s, x + Inches(0.22), Inches(4.65), Inches(3.55))
        add_text(s, x + Inches(0.22), Inches(4.85), Inches(3.55), Inches(0.85),
                 [[(fit, 13, INK, False)]], line_spacing=1.18)
        add_text(s, x + Inches(0.22), Inches(5.75), Inches(3.55), Inches(0.75),
                 [[(notfit, 12, GRAY, False)]], line_spacing=1.18)
    footer(s, idx)
    add_notes(s, "白板三件按任务拆：好看的看板、带洞察的图表、XRD 科研图。共用飞书介质，不共用工作流。")


def slide_feishu_whiteboard(prs, idx):
    s = blank(prs)
    kicker_bar(s, "11  细读", "beautiful-feishu-whiteboard",
               "github.com/zarazhangrui/beautiful-feishu-whiteboard  ·  681 Star  ·  MIT")
    facts = [
        ("35 套色板", "克制 / 均衡 / 大胆三档。同一套 LLM 训练三阶段内容，方便直接比气质。"),
        ("硬限制写进 RULES.md", "只用原生形状、无透明度、无渐变模糊；文字颜色导出有已知坑。全是实机踩出来的。"),
        ("模板不管构图", "templates/<slug>/design.md 只约束色板和情绪。布局由 Agent 按内容决定。"),
        ("交付物", "写入你自己的飞书租户：文档链接 + 预览图。随时可换风格重出。"),
        ("依赖", "Node 20+、飞书账号、全局 lark-cli 扫码登录、npx @larksuite/whiteboard-cli。"),
        ("安装", "npx skills add zarazhangrui/beautiful-feishu-whiteboard"),
    ]
    for i, (t, b) in enumerate(facts):
        col, row = i % 2, i // 2
        x = Inches(0.5) + col * Inches(6.35)
        y = Inches(1.45) + row * Inches(1.7)
        add_rect(s, x, y, Inches(6.15), Inches(1.55), WHITE)
        add_text(s, x + Inches(0.25), y + Inches(0.18), Inches(5.7), Inches(0.32),
                 [[(t, 15, INK, True)]])
        add_text(s, x + Inches(0.25), y + Inches(0.55), Inches(5.7), Inches(0.8),
                 [[(b, 13, MUTED, False)]], line_spacing=1.2)
    footer(s, idx)
    add_notes(s, "张咋啦白板：色板技能，不是图表技能。构图自由，颜色和介质限制不自由。")


def slide_graphy(prs, idx):
    s = blank(prs)
    kicker_bar(s, "12  细读", "larkboard-graphy  ·  数据叙事优先",
               "github.com/lostvita/larkboard-graphy  ·  14 Star  ·  MIT  ·  作者 F俊W")
    add_rect(s, Inches(0.5), Inches(1.45), Inches(12.35), Inches(1.15), INK)
    add_text(s, Inches(0.75), Inches(1.55), Inches(11.9), Inches(0.95),
             [[("普通图表回答「是多少」。这套技能逼你先回答「所以呢」。", 18, WHITE, True)],
              [("Widget 是一等公民：在画第一根柱子之前，先提出注释、高亮、趋势和对比策略。", 13, PAPER, False)]],
             line_spacing=1.2)
    left_items = [
        ("三段式", "标题区结论先行 · 图表区数据+Widget · 信息区洞察和来源。"),
        ("图表", "柱 / 堆叠 / 百分比堆叠 / 条 / 折线 / 饼 / 环 / 柱线组合。"),
        ("校验", "whiteboard-cli --check 查溢出和重叠，最多自动修 3 轮再写入飞书。"),
    ]
    for i, (k, v) in enumerate(left_items):
        y = Inches(2.8) + i * Inches(0.95)
        add_rect(s, Inches(0.5), y, Inches(6.15), Inches(0.85), WHITE)
        add_text(s, Inches(0.7), y + Inches(0.1), Inches(5.75), Inches(0.28),
                 [[(k, 13, RED, True)]])
        add_text(s, Inches(0.7), y + Inches(0.4), Inches(5.75), Inches(0.35),
                 [[(v, 12, MUTED, False)]])
    widgets = [
        ("Comment", "解释为什么"),
        ("PinNumber", "钉住关键数"),
        ("Sticker", "情绪标记"),
        ("DifferenceArrow", "量化差距"),
        ("AverageLine / GoalLine", "均值与目标"),
        ("TrendLine", "回归趋势"),
        ("Highlight", "其余淡到 0.3"),
        ("HighlightLabel", "只标重点"),
    ]
    add_rect(s, Inches(6.85), Inches(2.8), Inches(6.0), Inches(3.7), WHITE)
    add_text(s, Inches(7.05), Inches(2.95), Inches(5.6), Inches(0.3),
             [[("Widget 叙事层", 13, RED, True)]])
    for i, (n, d) in enumerate(widgets):
        y = Inches(3.35) + i * Inches(0.37)
        add_text(s, Inches(7.05), y, Inches(2.6), Inches(0.35),
                 [[(n, 12, INK, True)]], anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(9.7), y, Inches(2.9), Inches(0.35),
                 [[(d, 12, MUTED, False)]], anchor=MSO_ANCHOR.MIDDLE)
    footer(s, idx)
    add_notes(s, "larkboard-graphy 把图表做成汇报叙事。中文 README 里的 github.com/user/... 是占位，真实仓库是 lostvita/larkboard-graphy。")


def slide_xrd(prs, idx):
    s = blank(prs)
    kicker_bar(s, "13  细读", "XRD-SKILL  ·  科研绘图一图流",
               "github.com/LucianaiB2004/XRD-SKILL  ·  3 Star  ·  MIT  ·  聊天原文有空格，仓库名无空格")
    steps = [
        ("1", "读数", "二维 XRD 文本，默认匹配 XN*_Theta_2-Theta.txt。"),
        ("2", "出图", "堆叠谱线 SVG / PNG，可选 OpenAPI 节点 JSON。"),
        ("3", "上板", "写入已有白板，或在文档/Wiki 追加新画板块。"),
        ("4", "可改", "飞书里继续拖节点，不是一张死图。"),
    ]
    for i, (n, t, b) in enumerate(steps):
        x = Inches(0.5) + i * Inches(3.2)
        add_rect(s, x, Inches(1.45), Inches(3.05), Inches(2.15), WHITE)
        add_text(s, x + Inches(0.18), Inches(1.6), Inches(2.7), Inches(0.4),
                 [[(n, 22, RED, True)]])
        add_text(s, x + Inches(0.18), Inches(2.05), Inches(2.7), Inches(0.35),
                 [[(t, 16, INK, True)]])
        add_text(s, x + Inches(0.18), Inches(2.5), Inches(2.7), Inches(0.85),
                 [[(b, 13, MUTED, False)]], line_spacing=1.18)
    notes = [
        ("脚本主导", "xrd_data_to_chart.py 和 publish_xrd_whiteboard.py 固定流程，避免每次让模型重写画图胶水。"),
        ("环境", "Python 3.10+、Node 20+、npx @larksuite/whiteboard-cli、已认证 lark-cli。可 --skip-lark 只验本地出图。"),
        ("安装", "python scripts/install_codex_skill.py --force 后重启 Codex，口令 $xrd-onepage-whiteboard。"),
        ("边界", "主打 XRD 谱线一图流，不是通用 Origin 替代，也不是任意科研图模板库。"),
    ]
    for i, (t, b) in enumerate(notes):
        col, row = i % 2, i // 2
        x = Inches(0.5) + col * Inches(6.35)
        y = Inches(3.8) + row * Inches(1.5)
        add_rect(s, x, y, Inches(6.15), Inches(1.35), WHITE)
        add_text(s, x + Inches(0.25), y + Inches(0.15), Inches(5.7), Inches(0.3),
                 [[(t, 14, INK, True)]])
        add_text(s, x + Inches(0.25), y + Inches(0.5), Inches(5.7), Inches(0.7),
                 [[(b, 13, MUTED, False)]], line_spacing=1.18)
    footer(s, idx)
    add_notes(s, "XRD 一图流是清单里最窄、也最可复现的一项。聊天链接里的空格是笔误。")


def slide_scenes(prs, idx):
    s = blank(prs)
    kicker_bar(s, "14  选型", "十二件工具，十个场景")
    rows = [
        ("明天领导要改 PPT", "ppt-master / 花叔", "只要别人能在 WPS 里改，就走 PPTX。"),
        ("自己上台，要好看", "frontend-slides / 归藏", "双击浏览器。计时讲稿再加 Lewis。"),
        ("已有文案，只要气质", "beautiful-html-templates", "三封面预览，选定再填。"),
        ("HTML 演示要改两个字", "htmledit.ai", "点选即改，不必回 Agent 重生成。"),
        ("发群、做课件卡片", "baoyu-slide-deck", "整页图，为刷和转发。"),
        ("一堆公众号/PDF 先理清", "乔木", "结构初稿，再用上面任意一款做终稿。"),
        ("飞书里画架构/流程", "张咋啦白板", "35 套色板，原生可编辑。"),
        ("把Excel变成带洞察的汇报图", "larkboard-graphy", "先叙事策略，再画柱线环。"),
        ("XRD 谱线要交可改图", "XRD-SKILL", "原始数据到飞书画板一条链。"),
        ("品牌发布、动画、MP4", "huashu-design", "工作室流程，HTML 和 PPTX 都能要。"),
    ]
    for i, (scene, skill, why) in enumerate(rows):
        col, row = i % 2, i // 2
        x = Inches(0.5) + col * Inches(6.35)
        y = Inches(1.4) + row * Inches(1.05)
        add_rect(s, x, y, Inches(6.15), Inches(0.95), WHITE)
        add_rect(s, x, y, Inches(0.08), Inches(0.95), RED)
        add_text(s, x + Inches(0.25), y + Inches(0.08), Inches(5.7), Inches(0.28),
                 [[(scene, 13, INK, True)]])
        add_text(s, x + Inches(0.25), y + Inches(0.36), Inches(5.7), Inches(0.22),
                 [[(skill, 12, RED, True)]])
        add_text(s, x + Inches(0.25), y + Inches(0.58), Inches(5.7), Inches(0.28),
                 [[(why, 11, MUTED, False)]])
    footer(s, idx)
    add_notes(s, "选型页把四条赛道收成十个场景。能对上就用，对不上就不要硬塞。")


def slide_install(prs, idx):
    s = blank(prs)
    kicker_bar(s, "15  上手", "安装与入口",
               "演示七款用 npx skills add。模板库、改稿器、白板另有自己的入口。")
    cmds = [
        ("ppt-master", "npx skills add hugohe3/ppt-master"),
        ("frontend-slides", "npx skills add zarazhangrui/frontend-slides"),
        ("guizang-ppt-skill", "npx skills add op7418/guizang-ppt-skill"),
        ("huashu-design", "npx skills add alchaincyf/huashu-design"),
        ("html-ppt-skill", "npx skills add lewislulu/html-ppt-skill"),
        ("baoyu-slide-deck", "npx skills add JimLiu/baoyu-skills --skill baoyu-slide-deck"),
        ("qiaomu", "npx skills add joeseesun/qiaomu-anything-to-notebooklm"),
        ("html-templates", "克隆仓库，让 Agent 读 AGENTS.md"),
        ("feishu-whiteboard", "npx skills add zarazhangrui/beautiful-feishu-whiteboard"),
        ("larkboard-graphy", "git clone lostvita/larkboard-graphy → ~/.cursor/skills/"),
        ("htmledit", "打开 htmledit.ai  ·  或本地 python3 -m http.server"),
        ("XRD-SKILL", "python scripts/install_codex_skill.py --force"),
    ]
    for i, (name, cmd) in enumerate(cmds):
        col, row = i % 2, i // 2
        x = Inches(0.45) + col * Inches(6.4)
        y = Inches(1.4) + row * Inches(0.88)
        add_rect(s, x, y, Inches(6.2), Inches(0.78), WHITE)
        add_text(s, x + Inches(0.15), y + Inches(0.06), Inches(5.9), Inches(0.28),
                 [[(name, 12, RED, True)]])
        add_text(s, x + Inches(0.15), y + Inches(0.36), Inches(5.9), Inches(0.35),
                 [[(cmd, 11, INK, False)]])
    footer(s, idx)
    add_notes(s, "白板三件都要飞书登录。财猫非商业协议。归藏是 AGPL-3.0。")


def slide_judgements(prs, idx):
    s = blank(prs)
    kicker_bar(s, "16  判断", "六句把这张地图读完")
    items = [
        ("01", "七款 PPT 是赛道 A，不是整张地图。",
         "许俊贤补上的是生成之后：改 HTML、上白板、画科研图。只盯着 Star 会漏掉真正缺的那一环。"),
        ("02", "张咋啦三件套是一条生产线。",
         "模板库选气质，frontend-slides 做演示，飞书白板做可编辑看板。不要三选一，按介质串联。"),
        ("03", "HTML 路线缺的不是生成，是改。",
         "htmledit.ai 让不会写代码的人能改 frontend-slides / 归藏 / Lewis 的产物。这是演示赛道的下游。"),
        ("04", "飞书白板三件不要互相替代。",
         "色板看板、叙事图表、XRD 一图流共用介质，工作流完全不同。选错会做出一张好看但没洞察的图。"),
        ("05", "协议和账号是隐形成本。",
         "归藏 AGPL；财猫非商业；白板三件都要飞书租户和 CLI 登录。能跑通环境，才谈得上「装了」。"),
        ("06", "选工具之前写一句：这份东西第二天交给谁、在哪个软件里打开。",
         "WPS 改稿、浏览器开讲、飞书里继续画、实验室交谱图——这句话比任何清单都管用。"),
    ]
    for i, (num, title, body) in enumerate(items):
        y = Inches(1.38) + i * Inches(0.9)
        add_text(s, Inches(0.5), y, Inches(0.7), Inches(0.8),
                 [[(num, 16, RED, True)]])
        add_text(s, Inches(1.2), y, Inches(11.6), Inches(0.32),
                 [[(title, 14, INK, True)]])
        add_text(s, Inches(1.2), y + Inches(0.34), Inches(11.6), Inches(0.48),
                 [[(body, 12, MUTED, False)]])
    footer(s, idx)
    add_notes(s, "六条判断是结论。如果只能留一页，留这一页。")


def slide_relation(prs, idx):
    s = blank(prs)
    kicker_bar(s, "17  关系", "和上一份七款横评怎么一起用")
    cols = [
        ("上一份：七款横评",
         "问的是「哪款能做出像样的 PPT」。",
         ["社区锐评排名", "三条产物路线：PPTX / HTML / 图片", "七款各一页细读", "办公流转 vs 现场演示"]),
        ("这一份：技能清单",
         "问的是「许俊贤甩过来的地图有多宽」。",
         ["不排名，按赛道", "补上模板库、改稿器、白板", "张咋啦从 1 件变成 3 件", "科研 XRD 单独成环"]),
        ("建议用法",
         "先用这份找赛道，再用横评做演示细选。",
         ["演示需求 → 打开横评第 14–16 页", "改 HTML → 直接 htmledit.ai", "飞书看板/图表/XRD → 本清单赛道 D", "不要用 Star 给白板和 PPT 排名"]),
    ]
    for i, (title, sub, bullets) in enumerate(cols):
        x = Inches(0.5) + i * Inches(4.2)
        add_rect(s, x, Inches(1.45), Inches(4.0), Inches(5.4), WHITE)
        add_rect(s, x, Inches(1.45), Inches(4.0), Inches(0.08), RED)
        add_text(s, x + Inches(0.22), Inches(1.7), Inches(3.55), Inches(0.7),
                 [[(title, 16, INK, True)]], line_spacing=1.05)
        add_text(s, x + Inches(0.22), Inches(2.45), Inches(3.55), Inches(0.7),
                 [[(sub, 13, MUTED, False)]], line_spacing=1.18)
        hairline(s, x + Inches(0.22), Inches(3.25), Inches(3.55))
        y = Inches(3.45)
        for b in bullets:
            add_rect(s, x + Inches(0.22), y + Inches(0.08), Inches(0.08), Inches(0.08), RED)
            add_text(s, x + Inches(0.45), y, Inches(3.3), Inches(0.7),
                     [[(b, 13, INK, False)]])
            y += Inches(0.8)
    footer(s, idx)
    add_notes(s, "两份材料互补：横评解决演示赛道细选，清单解决地图宽度。")


def slide_close(prs, idx):
    s = blank(prs)
    set_bg(s, INK)
    add_rect(s, 0, 0, Inches(0.18), PAGE_H, RED)
    add_text(s, Inches(0.7), Inches(1.5), Inches(12), Inches(0.4),
             [[("十二个名字可以忘掉", 14, RED, True)]])
    add_text(s, Inches(0.7), Inches(2.05), Inches(12), Inches(1.6),
             [[("记住四件事", 36, WHITE, True)],
              [("演示、模板、改稿、白板", 36, WHITE, True)]],
             line_spacing=1.05)
    add_rect(s, Inches(0.7), Inches(4.0), Inches(1.6), Inches(0.06), RED)
    add_text(s, Inches(0.7), Inches(4.3), Inches(11.8), Inches(1.4),
             [[("路演交 ppt-master。上台交张咋啦或归藏。改两个字打开 htmledit.ai。", 16, PAPER, False)],
              [("飞书里画看板用色板，画汇报图用 graphy，交 XRD 用一图流。", 16, PAPER, False)]],
             line_spacing=1.28)
    add_text(s, Inches(0.7), Inches(6.2), Inches(11.8), Inches(0.5),
             [[("资料截止 2026-08-23  ·  链接以各细读页为准  ·  可在 PowerPoint / WPS 中继续编辑", 12, FOOT_GRAY, False)]])
    add_notes(s, "收束：把十二个名字还回给四条赛道。问听众下一次要交的文件在哪个软件里打开。")


def build():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = PAGE_W
    prs.slide_height = PAGE_H

    slide_cover(prs)
    slide_source(prs, 2)
    slide_how(prs, 3)
    slide_routes(prs, 4)
    slide_inventory(prs, 5)
    slide_deck_grid(prs, 6)
    slide_deck_pick(prs, 7)
    slide_zhang_family(prs, 8)
    slide_html_templates(prs, 9)
    slide_htmledit(prs, 10)
    slide_board_overview(prs, 11)
    slide_feishu_whiteboard(prs, 12)
    slide_graphy(prs, 13)
    slide_xrd(prs, 14)
    slide_scenes(prs, 15)
    slide_install(prs, 16)
    slide_judgements(prs, 17)
    slide_relation(prs, 18)
    slide_close(prs, 19)

    actual = len(prs.slides)
    if actual != TOTAL:
        print(f"警告：预期 {TOTAL} 页，实际 {actual} 页")
    prs.save(str(OUT))
    size_kb = OUT.stat().st_size / 1024
    print(f"已生成：{OUT}")
    print(f"页数：{actual}    大小：{size_kb:.1f} KB")
    return str(OUT), actual


if __name__ == "__main__":
    build()
