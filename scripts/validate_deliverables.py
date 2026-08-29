#!/usr/bin/env python3
"""校验三份交付物结构与关键结论是否写入。"""

from pathlib import Path
import zipfile
import xml.etree.ElementTree as ET

from openpyxl import load_workbook
from pptx import Presentation
from docx import Document

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "deliverables"

MUST = [
    "胡**",
    "64",
    "胫骨远端",
    "腓骨近端",
    "半脱位",
    "伤残尚未",
    "内外固定",
    "跟骨骨刺",
    "人民医院",
    "保险公司",
    "美团",
    "康复期",
    "原医院结算",
]


def docx_text(path: Path) -> str:
    doc = Document(path)
    parts = [p.text for p in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                parts.append(cell.text)
    return "\n".join(parts)


def ppt_text(path: Path) -> str:
    prs = Presentation(path)
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.append(shape.text_frame.text)
    return "\n".join(chunks)


def pdf_text(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        return ""
    reader = PdfReader(str(path))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def main() -> None:
    docx = OUT / "青岛抚顺路和哈尔滨路路口交通事故_伤情与处理备忘录_20260817.docx"
    xlsx = OUT / "青岛抚顺路和哈尔滨路路口交通事故_伤情伤残与行动表_20260817.xlsx"
    pptx = OUT / "青岛抚顺路和哈尔滨路路口交通事故_伤情与伤残简报_20260817.pptx"
    comm_docx = OUT / "青岛抚顺路和哈尔滨路路口交通事故_肇事方沟通方案与体检分析_20260817.docx"
    comm_pdf = OUT / "青岛抚顺路和哈尔滨路路口交通事故_肇事方沟通方案与体检分析_20260817.pdf"
    comm_xlsx = OUT / "青岛抚顺路和哈尔滨路路口交通事故_肇事方沟通流程表_20260817.xlsx"
    for p in (docx, xlsx, pptx, comm_docx, comm_pdf, comm_xlsx):
        assert p.exists() and p.stat().st_size > 2000, p

    text = docx_text(docx)
    print("DOCX 段落+表格字符", len(text), "文件字节", docx.stat().st_size)
    for k in MUST:
        assert k in text, f"Word 缺少：{k}"
    assert "刘孝春" in text
    assert "补差额" in text
    assert "尚未出具" in text
    assert "乔刘记商贸" in text
    assert "驾驶资格" in text
    assert "不是借款" in text
    assert "胡志远" in text
    assert "0001519455" in text
    assert "开放性" in text
    assert "VSD" in text
    assert "总赔偿数额" in text
    assert "雇员受害" in text
    assert "不要去人社局" in text or "不去人社局" in text
    assert "逐笔核清" in text or "4 万元" in text
    assert "原医院结算" in text
    assert "工伤或用工" in text
    assert "立即八事" in text

    wb = load_workbook(xlsx)
    sheets = wb.sheetnames
    print("XLSX 工作表", sheets)
    need = ["办理顺序", "材料核心", "伤情鉴定", "伤残情况", "外伤与退变对照", "转院报销", "费用台账", "72小时待办", "沟通口径卡", "责任与程序", "刘孝春分担清单", "法律关系提醒", "待补充信息"]
    for n in need:
        assert n in sheets, n
    stage_ws = wb["办理顺序"]
    stage_text = "\n".join(str(c.value or "") for row in stage_ws.iter_rows(max_row=16) for c in row)
    assert "总赔偿数额" in stage_text
    assert "监控" in stage_text and "用工" in stage_text
    core_ws = wb["材料核心"]
    core_sheet = "\n".join(str(c.value or "") for row in core_ws.iter_rows(max_row=20) for c in row)
    assert "原医院结算" in core_sheet and "4 万元" in core_sheet
    todo_ws = wb["72小时待办"]
    todo_text = "\n".join(str(c.value or "") for row in todo_ws.iter_rows(max_row=12) for c in row)
    assert "4 万元" in todo_text and "书面事故认定书" in todo_text
    ws = wb["伤情鉴定"]
    joined = "\n".join(str(c.value or "") for row in ws.iter_rows(max_row=8) for c in row)
    assert "10008056847" in joined and "10008059257" in joined
    dis = wb["伤残情况"]
    dis_text = "\n".join(str(c.value or "") for row in dis.iter_rows(max_row=10) for c in row)
    assert "尚未" in dis_text or "不能" in dis_text

    prs = Presentation(pptx)
    print("PPT 页数", len(prs.slides))
    assert len(prs.slides) == 10
    ptext = ppt_text(pptx)
    for k in ("伤残尚未鉴定", "胫骨远端", "禁止说"):
        assert k in ptext.replace("\n", "") or k in ptext, k

    ctext = docx_text(comm_docx)
    print("沟通方案 DOCX 字符", len(ctext), "字节", comm_docx.stat().st_size)
    for k in ("胫骨远端", "腓骨近端", "跟骨骨刺", "开场", "律师函", "方案甲", "12", "人民医院", "美团"):
        assert k in ctext, f"沟通方案 Word 缺少：{k}"
    assert "不要请诉讼律师" in ctext or "现在不要请诉讼律师" in ctext or "不到场见骑手" in ctext or "不先发律师函" in ctext

    cwb = load_workbook(comm_xlsx)
    print("沟通流程 XLSX 工作表", cwb.sheetnames)
    for n in ("怎么用", "12步沟通顺序", "对方一句话怎么接", "律师决策", "三套方案", "影像三份对照"):
        assert n in cwb.sheetnames, n
    flow = cwb["12步沟通顺序"]
    flow_text = "\n".join(str(c.value or "") for row in flow.iter_rows(max_row=16) for c in row)
    assert "未开始" in flow_text and "肇事女骑手" in flow_text

    boss_ws = wb["刘孝春分担清单"]
    boss_text = "\n".join(str(c.value or "") for row in boss_ws.iter_rows(max_row=20) for c in row)
    assert "护理费" in boss_text and "误工费" in boss_text and "未结工钱" in boss_text

    ptext_c = pdf_text(comm_pdf)
    print("沟通方案 PDF 抽取字符", len(ptext_c), "字节", comm_pdf.stat().st_size)
    assert comm_pdf.stat().st_size > 8000
    if ptext_c:
        for k in ("胫骨远端", "方案甲", "律师", "开场", "10008056847", "胡继刚"):
            assert k in ptext_c.replace(" ", ""), f"沟通方案 PDF 缺少：{k}"

    html = OUT / "hongfeng-guide.html"
    html_cn = OUT / "青岛抚顺路和哈尔滨路路口交通事故_处理总览.html"
    for p in (html, html_cn):
        assert p.exists() and p.stat().st_size > 8000, p
        htext = p.read_text(encoding="utf-8")
        assert "<script" not in htext
        assert "红枫路" not in htext, f"网页仍有红枫路：{p}"
        for k in ("青岛抚顺路和哈尔滨路路口", "胫骨远端", "人民医院", "10008056847", "美团", "刘孝春", "护理费", "还缺", "护栏开口", "乔刘记", "尚未出具", "总赔偿数额", "4 万元"):
            assert k in htext, f"网页缺少：{k}"
    print("HTML 字节", html.stat().st_size)

    asr_md = OUT / "交通事故材料整理_2026-08-19.md"
    assert asr_md.exists() and asr_md.stat().st_size > 8000, asr_md
    asr_text = asr_md.read_text(encoding="utf-8")
    for k in ("护栏开口", "录音转写", "刘小春", "九级", "4 万元", "原医院结算"):
        assert k in asr_text, f"8/19 整理稿缺少：{k}"

    corpus = OUT / "沟通记录合集-2026年8月15日至18日.md"
    guide = OUT / "沟通记录合集_口径对照_2026-08-19.md"
    assert corpus.exists() and corpus.stat().st_size > 100_000, corpus
    assert guide.exists() and guide.stat().st_size > 3000, guide
    gtext = guide.read_text(encoding="utf-8")
    for k in ("同一起", "远洋", "33 床", "1 号门", "踏板", "1866"):
        assert k in gtext, f"口径对照缺少：{k}"
    assert "拆成多起事故是错的" in corpus.read_text(encoding="utf-8")[:2000]

    ans = OUT / "待确认事项_家属答复_2026-08-19.md"
    assert ans.exists() and ans.stat().st_size > 2000, ans
    atxt = ans.read_text(encoding="utf-8")
    for k in ("尚未出具", "乔刘记商贸", "不是借款", "抚顺路和哈尔滨路路口", "胡志远病例文件", "19:01"):
        assert k in atxt, f"家属答复缺少：{k}"

    core_md = OUT / "材料核心归结_立即八事.md"
    assert core_md.exists() and core_md.stat().st_size > 1500, core_md
    core_text = core_md.read_text(encoding="utf-8")
    for k in ("原医院结算", "书面责任认定", "4 万元", "用工责任", "伤残鉴定", "工伤或用工"):
        assert k in core_text, f"材料核心归结缺少：{k}"

    stage_md = OUT / "现阶段办理顺序_紧急程度.md"
    assert stage_md.exists() and stage_md.stat().st_size > 3000, stage_md
    stage_md_text = stage_md.read_text(encoding="utf-8")
    for k in ("总赔偿数额", "紧急", "书面责任认定", "人社局", "雇员受害", "怎么办理"):
        assert k in stage_md_text, f"办理顺序手册缺少：{k}"

    lawyer_md = OUT / "给律师的事故完整经过说明_2026-08-19.md"
    lawyer_docx = OUT / "青岛抚顺路和哈尔滨路路口交通事故_给律师的完整经过说明_20260819.docx"
    lawyer_pdf = OUT / "青岛抚顺路和哈尔滨路路口交通事故_给律师的完整经过说明_20260819.pdf"
    assert lawyer_md.exists() and lawyer_md.stat().st_size > 3000, lawyer_md
    ltxt = lawyer_md.read_text(encoding="utf-8")
    for k in ("完整经过", "抚顺路和哈尔滨路路口", "抚顺路批发市场", "脚踏板", "10008056847", "尚未出具", "乔刘记商贸", "4 万元", "不是律师函", "胡志远", "0001519455", "开放性", "VSD", "病历图文"):
        assert k in ltxt, f"律师经过说明缺少：{k}"
    assert "红枫路" not in ltxt, "律师经过说明不得再出现红枫路"
    assert lawyer_docx.exists() and lawyer_docx.stat().st_size > 4000, lawyer_docx
    ldoc = docx_text(lawyer_docx)
    for k in ("完整经过", "内外固定", "刘孝春", "简易", "民法典", "抚顺路和哈尔滨路路口", "胡志远", "19:01"):
        assert k in ldoc, f"律师经过 Word 缺少：{k}"
    assert "红枫路" not in ldoc, "律师经过 Word 不得再出现红枫路"
    assert lawyer_pdf.exists() and lawyer_pdf.stat().st_size > 8000, lawyer_pdf
    lp = pdf_text(lawyer_pdf)
    if lp:
        for k in ("完整经过", "抚顺路和哈尔滨路路口", "刘孝春", "内外固定", "胡志远"):
            assert k in lp.replace(" ", ""), f"律师经过 PDF 缺少：{k}"
        assert "红枫路" not in lp, "律师经过 PDF 不得再出现红枫路"

    chart_pdf = OUT / "胡志远病例文件.pdf"
    chart_md = OUT / "齐鲁医院_胡志远病例摘录_2026-08-19.md"
    assert chart_pdf.exists() and chart_pdf.stat().st_size > 100_000, chart_pdf
    assert chart_md.exists() and "0001519455" in chart_md.read_text(encoding="utf-8")

    illu_docx = OUT / "青岛抚顺路和哈尔滨路路口交通事故_给律师的病历图文_20260819.docx"
    illu_pdf = OUT / "青岛抚顺路和哈尔滨路路口交通事故_给律师的病历图文_20260819.pdf"
    photo_dir = OUT / "chart-photos"
    photos = sorted(photo_dir.glob("*.jpg"))
    assert illu_docx.exists() and illu_docx.stat().st_size > 200_000, illu_docx
    assert illu_pdf.exists() and illu_pdf.stat().st_size > 200_000, illu_pdf
    assert len(photos) >= 7, f"chart-photos 只有 {len(photos)} 张"
    itxt = docx_text(illu_docx)
    for k in ("胡志远", "出院记录", "不是律师函", "开放性", "抚顺路和哈尔滨路路口", "19:01", "VSD"):
        assert k in itxt, f"病历图文 Word 缺少：{k}"
    assert "红枫路" not in itxt, "病历图文 Word 不得再出现红枫路"
    with zipfile.ZipFile(illu_docx) as zf:
        media = [n for n in zf.namelist() if n.startswith("word/media/")]
        assert len(media) >= 7, f"Word 内嵌图片只有 {len(media)} 张"
    print("病历图文 DOCX 字节", illu_docx.stat().st_size, "PDF 字节", illu_pdf.stat().st_size, "抽出照片", len(photos))

    kimi = OUT / "事故3D复原_Kimi提示词.md"
    kimi_text = kimi.read_text(encoding="utf-8")
    assert "待用" in kimi_text and "护栏开口" in kimi_text
    assert "对向" in kimi_text and "我不是闪的" in kimi_text
    seedance = OUT / "事故现场复盘_Seedance2.5提示词.md"
    assert seedance.exists() and seedance.stat().st_size > 3000, seedance
    stext = seedance.read_text(encoding="utf-8")
    for k in ("Seedance 2.5", "待用", "抚顺路和哈尔滨路路口", "护栏开口", "脚踏板", "不是律师函", "非原始视频", "1 号门", "上面车也下来了", "我不是闪的", "对向", "9:16", "480P", "42"):
        assert k in stext, f"Seedance 提示词缺少：{k}"
    assert "红枫路" not in stext
    mp4 = OUT / "事故3D复原_示范动画.mp4"
    assert mp4.exists() and mp4.stat().st_size > 50_000, mp4

    liu_md = OUT / "给刘孝春的垫付欠薪与护理说明_2026-08-19.md"
    liu_docx = OUT / "青岛抚顺路和哈尔滨路路口交通事故_给刘孝春的垫付与护理说明_20260819.docx"
    liu_pdf = OUT / "青岛抚顺路和哈尔滨路路口交通事故_给刘孝春的垫付与护理说明_20260819.pdf"
    assert liu_md.exists() and liu_md.stat().st_size > 2000, liu_md
    lmd = liu_md.read_text(encoding="utf-8")
    for k in ("不是借款", "不是了结", "一对一", "不是律师函", "抚顺路和哈尔滨路路口", "不冲抵", "不代美团"):
        assert k in lmd, f"给刘说明 Markdown 缺少：{k}"
    assert "家属自留" in lmd and "不要打印给刘" in lmd
    assert "红枫路" not in lmd
    assert liu_docx.exists() and liu_docx.stat().st_size > 4000, liu_docx
    ldoc = docx_text(liu_docx)
    for k in ("不是借款", "不是了结", "一对一", "不是律师函", "抚顺路和哈尔滨路路口", "不冲抵未结工钱", "不代美团"):
        assert k in ldoc, f"给刘说明 Word 缺少：{k}"
    assert "红枫路" not in ldoc
    assert "30% 全包" not in ldoc and "30%全包" not in ldoc
    assert "70%" not in ldoc and "30%" not in ldoc
    assert "家属自留" not in ldoc
    assert liu_pdf.exists() and liu_pdf.stat().st_size > 8000, liu_pdf
    lp_liu = pdf_text(liu_pdf)
    if lp_liu:
        compact = lp_liu.replace(" ", "")
        for k in ("不是借款", "一对一", "不是律师函", "抚顺路和哈尔滨路路口"):
            assert k in compact, f"给刘说明 PDF 缺少：{k}"
        assert "红枫路" not in lp_liu
        assert "30%全包" not in compact and "家属自留" not in lp_liu

    mom_md = OUT / "给妈妈跟刘孝春沟通的口径_2026-08-20.md"
    mom_docx = OUT / "青岛抚顺路和哈尔滨路路口交通事故_给妈妈跟刘孝春的接话卡_20260820.docx"
    mom_pdf = OUT / "青岛抚顺路和哈尔滨路路口交通事故_给妈妈跟刘孝春的接话卡_20260820.pdf"
    assert mom_md.exists() and mom_md.stat().st_size > 2500, mom_md
    mtxt = mom_md.read_text(encoding="utf-8")
    for k in ("家属自留", "不是借款", "不是了结", "一对一", "开场", "不要发给刘孝春", "抚顺路和哈尔滨路路口"):
        assert k in mtxt, f"妈妈口径 Markdown 缺少：{k}"
    assert "红枫路" not in mtxt
    assert mom_docx.exists() and mom_docx.stat().st_size > 4000, mom_docx
    mdoc = docx_text(mom_docx)
    for k in ("不是借款", "不是了结", "一对一", "家属自留", "不要给刘孝春", "开场"):
        assert k in mdoc, f"妈妈接话卡 Word 缺少：{k}"
    assert "红枫路" not in mdoc
    assert "30% 全包" not in mdoc and "30%全包" not in mdoc
    assert "70%" not in mdoc and "30%" not in mdoc
    assert "民法典" not in mdoc
    assert mom_pdf.exists() and mom_pdf.stat().st_size > 8000, mom_pdf
    mp = pdf_text(mom_pdf)
    if mp:
        compact_m = mp.replace(" ", "")
        for k in ("不是借款", "一对一", "家属自留", "不要给刘孝春"):
            assert k in compact_m, f"妈妈接话卡 PDF 缺少：{k}"
        assert "红枫路" not in mp
        assert "70%" not in mp and "30%" not in mp

    rider_md = OUT / "给美团骑手的跟进微信_2026-08-22.md"
    assert rider_md.exists() and rider_md.stat().st_size > 800, rider_md
    rtxt = rider_md.read_text(encoding="utf-8")
    for k in ("拉群", "不是让你个人现在拿钱", "市北区人民医院", "不是律师函", "今天不谈一共赔多少"):
        assert k in rtxt, f"骑手跟进微信缺少：{k}"
    assert "红枫路" not in rtxt
    copy = rtxt.split("## 请直接复制发出")[1].split("## 她若回了")[0]
    assert "全责" not in copy
    assert "律师函" not in copy
    assert "不先发律师函" in rtxt

    itinerary = OUT / "胡继刚行程_长春上海青岛_2026-08-22.md"
    assert itinerary.exists() and itinerary.stat().st_size > 1200, itinerary
    itxt = itinerary.read_text(encoding="utf-8")
    for k in ("长春", "一对一", "护理费", "误工费", "杨浦", "家属自留", "不要发给刘孝春"):
        assert k in itxt, f"行程说明缺少：{k}"
    assert "红枫路" not in itxt

    confirm = OUT / "家属确认_岳父律师与护工_2026-08-22.md"
    assert confirm.exists() and confirm.stat().st_size > 1200, confirm
    ctxt = confirm.read_text(encoding="utf-8")
    for k in ("三十余年", "场外", "300", "一对一", "发票", "不是律师函", "不要发给"):
        assert k in ctxt, f"岳父护工确认缺少：{k}"
    assert "红枫路" not in ctxt
    assert "不先发" in ctxt or "不署名发律师函" in ctxt

    digest = OUT / "全案整理_截至2026-08-22.md"
    assert digest.exists() and digest.stat().st_size > 3000, digest
    dtxt = digest.read_text(encoding="utf-8")
    for k in ("抚顺路和哈尔滨路路口", "0001519455", "19:01", "尚未出具", "五句话", "300 元/天", "拉群", "再推迟", "杨浦", "不是律师函", "三十余年"):
        assert k in dtxt, f"全案整理缺少：{k}"
    assert "红枫路" not in dtxt

    pingan = OUT / "给平安理赔专员陈老师的回复_2026-08-29.md"
    assert pingan.exists() and pingan.stat().st_size > 2000, pingan
    ptxt_pa = pingan.read_text(encoding="utf-8")
    for k in ("陈老师", "报案号", "市北区人民医院", "护理费", "不报总数", "泰康不向平安披露", "认定书尚未出具", "分割单"):
        assert k in ptxt_pa, f"平安回复缺少：{k}"
    assert "红枫路" not in ptxt_pa

    sol_docx = OUT / "青岛抚顺路和哈尔滨路路口交通事故_完整解决方案_20260822.docx"
    sol_xlsx = OUT / "青岛抚顺路和哈尔滨路路口交通事故_赔付测算表_20260822.xlsx"
    assert sol_docx.exists() and sol_docx.stat().st_size > 10000, sol_docx
    stext_sol = docx_text(sol_docx)
    for k in ("泰康", "残疾赔偿金", "二次住院", "欠薪", "理赔", "家属内部", "不是律师函", "71703"):
        assert k in stext_sol, f"解决方案 Word 缺少：{k}"
    assert "红枫路" not in stext_sol
    swb = load_workbook(sol_xlsx)
    for n in ("责任与赔付分担", "理赔流程时间线", "住院与治疗时长", "费用估算明细", "三情景总账", "泰康保单核对", "纪律"):
        assert n in swb.sheetnames, n
    tot_ws = swb["三情景总账"]
    tot_text = "\n".join(str(c.value or "") for row in tot_ws.iter_rows(max_row=10) for c in row)
    assert "十级" in tot_text and "九级" in tot_text and "欠薪" in tot_text

    family = OUT / "家庭执行方案_2026-08-29.md"
    assert family.exists() and family.stat().st_size > 1500, family
    ftxt = family.read_text(encoding="utf-8")
    for k in ("托班", "求职", "二次住院", "授权委托书", "再推迟", "家属内部"):
        assert k in ftxt, f"家庭执行方案缺少：{k}"
    assert "红枫路" not in ftxt

    poa = OUT / "授权委托书_胡志远委托胡继刚_草稿.md"
    assert poa.exists() and poa.stat().st_size > 400, poa
    ptxt_poa = poa.read_text(encoding="utf-8")
    for k in ("胡志远", "胡继刚", "抚顺路和哈尔滨路路口", "事故认定书"):
        assert k in ptxt_poa, f"授权委托书缺少：{k}"
    assert "红枫路" not in ptxt_poa

    diet_md = OUT / "给护工的无糖床头食谱_2026-08-29.md"
    diet_docx = OUT / "青岛抚顺路和哈尔滨路路口交通事故_给护工的无糖床头食谱_20260829.docx"
    diet_pdf = OUT / "青岛抚顺路和哈尔滨路路口交通事故_给护工的无糖床头食谱_20260829.pdf"
    assert diet_md.exists() and diet_md.stat().st_size > 800, diet_md
    dtxt_diet = diet_md.read_text(encoding="utf-8")
    for k in ("早餐", "午餐", "晚餐", "加餐", "无糖", "低血糖", "胡志远", "2000"):
        assert k in dtxt_diet, f"床头食谱缺少：{k}"
    assert "红枫路" not in dtxt_diet
    assert diet_docx.exists() and diet_docx.stat().st_size > 5000, diet_docx
    assert diet_pdf.exists() and diet_pdf.stat().st_size > 3000, diet_pdf
    dtext_doc = docx_text(diet_docx)
    for k in ("早餐", "午餐", "晚餐", "禁止", "葡萄糖"):
        assert k in dtext_doc, f"床头食谱 Word 缺少：{k}"
    dtext_pdf = pdf_text(diet_pdf)
    if dtext_pdf:
        for k in ("早餐", "低血糖", "护工"):
            assert k in dtext_pdf, f"床头食谱 PDF 缺少：{k}"

    from pack_all import ZIP_NAME

    archive = OUT / ZIP_NAME
    assert archive.exists() and archive.stat().st_size > 1_000_000, archive
    with zipfile.ZipFile(archive) as zf:
        names = set(zf.namelist())
    for must in (
        "胡志远病例文件.pdf",
        "chart-photos/01_住院腕带.jpg",
        "青岛抚顺路和哈尔滨路路口交通事故_完整解决方案_20260822.docx",
        "家庭执行方案_2026-08-29.md",
        "授权委托书_胡志远委托胡继刚_草稿.md",
        "给护工的无糖床头食谱_2026-08-29.md",
        "青岛抚顺路和哈尔滨路路口交通事故_给护工的无糖床头食谱_20260829.pdf",
        "压缩包目录说明.txt",
    ):
        assert must in names, f"压缩包缺少：{must}"
    assert ZIP_NAME not in names

    assert "护栏开口" in text
    assert "待用" in text or "护栏" in text

    print("校验通过")


if __name__ == "__main__":
    main()
