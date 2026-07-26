# -*- coding: utf-8 -*-
"""
生成《复旦科技园创业孵化基地 2026 年度工作报告》领导汇报 PPT。

内容依据：source_docs/复旦科技园创业孵化基地2026年度工作报告.docx
版式与视觉参考：2023 年上海市创业孵化示范基地工作报告（复旦科技园）
素材：从原 PDF 提取的园区建筑、活动现场照片及 FDUSP Logo

运行：python3 scripts/build_2026_annual_report_ppt.py
输出：dist/复旦科技园创业孵化基地_2026年度工作报告.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "assets"
PHOTOS = ASSETS / "photos"
OUT = ROOT / "dist" / "复旦科技园创业孵化基地_2026年度工作报告.pptx"

# ----------------------------------------------------------------------------
# 视觉主题：深蓝政务风 + 复旦红点缀（对标 2023 汇报稿，面向领导汇报）
# ----------------------------------------------------------------------------
NAVY = RGBColor(0x0B, 0x2A, 0x4A)
NAVY_2 = RGBColor(0x12, 0x3A, 0x63)
NAVY_3 = RGBColor(0x1A, 0x4F, 0x7A)
STEEL = RGBColor(0x2F, 0x5F, 0x8A)
RED = RGBColor(0xC8, 0x10, 0x2E)
RED_SOFT = RGBColor(0xE8, 0x4A, 0x4A)
GOLD = RGBColor(0xC4, 0x9A, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF = RGBColor(0xF3, 0xF5, 0xF8)
LIGHT = RGBColor(0xE8, 0xED, 0xF3)
GREY = RGBColor(0x6B, 0x7A, 0x8D)
MUTED = RGBColor(0x8A, 0x97, 0xA8)
DARK = RGBColor(0x1A, 0x2A, 0x3A)
LINE = RGBColor(0xD0, 0xD7, 0xE0)

FONT = "微软雅黑"
FONT_FALLBACK = "WenQuanYi Micro Hei"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ----------------------------------------------------------------------------
# 基础绘图工具
# ----------------------------------------------------------------------------
def _set_font(run, size, color, bold=False, font=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", font)


def add_slide(bg=OFF):
    slide = prs.slides.add_slide(BLANK)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg
    rect.line.fill.background()
    rect.shadow.inherit = False
    return slide


def add_rect(slide, x, y, w, h, color, line_color=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line_color is not None:
        sp.line.color.rgb = line_color
        sp.line.width = Pt(line_w)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def add_text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: (text, size, color, bold[, space_after[, line_spacing]])"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    for i, spec in enumerate(lines):
        text, size, color, bold = spec[0], spec[1], spec[2], spec[3]
        space_after = spec[4] if len(spec) > 4 else 4
        line_spacing = spec[5] if len(spec) > 5 else None
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        _set_font(run, size, color, bold)
    return tb


def add_pic(slide, path, x, y, w=None, h=None):
    path = Path(path)
    if not path.exists():
        return None
    kwargs = {"left": x, "top": y}
    if w is not None:
        kwargs["width"] = w
    if h is not None:
        kwargs["height"] = h
    return slide.shapes.add_picture(str(path), **kwargs)


def header_bar(slide, section_no, title, light=True):
    """内容页顶栏：Logo + 章节标题"""
    add_rect(slide, 0, 0, SW, Inches(0.78), NAVY)
    # 左侧白底放置 Logo（对标 2023 汇报稿顶栏）
    add_rect(slide, 0, 0, Inches(2.35), Inches(0.78), WHITE)

    logo = ASSETS / "logo_stack.png"
    if logo.exists():
        add_pic(slide, logo, Inches(0.18), Inches(0.06), h=Inches(0.66))
    else:
        add_pic(slide, ASSETS / "logo_red_clean.png", Inches(0.2), Inches(0.12), h=Inches(0.52))

    add_text(
        slide,
        Inches(2.6),
        Inches(0.18),
        Inches(10.2),
        Inches(0.5),
        [(f"{section_no}  {title}", 18, WHITE, True)],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    # 红色细线
    add_rect(slide, 0, Inches(0.78), SW, Inches(0.035), RED)


def footer(slide, page_hint=""):
    add_rect(slide, 0, Inches(7.18), SW, Inches(0.32), NAVY)
    add_text(
        slide,
        Inches(0.45),
        Inches(7.2),
        Inches(9.5),
        Inches(0.28),
        [("复旦大学国家大学科技园 · 创业孵化基地  |  2026 年度工作报告", 9, MUTED, False)],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    if page_hint:
        add_text(
            slide,
            Inches(10.8),
            Inches(7.2),
            Inches(2.1),
            Inches(0.28),
            [(page_hint, 9, MUTED, False)],
            align=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE,
        )


def kpi_card(slide, x, y, w, h, value, label, sub=None, dark=False):
    bg = NAVY_2 if dark else WHITE
    vc = GOLD if dark else RED
    lc = WHITE if dark else DARK
    sc = MUTED if dark else GREY
    card = add_rect(slide, x, y, w, h, bg, line_color=None if dark else LINE, line_w=0.75, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    card.adjustments[0] = 0.08
    add_rect(slide, x, y, Inches(0.08), h, RED)
    add_text(slide, Emu(int(x) + Inches(0.22)), Emu(int(y) + Inches(0.18)), Emu(int(w) - Inches(0.35)), Inches(0.55),
             [(value, 26, vc, True)], align=PP_ALIGN.LEFT)
    add_text(slide, Emu(int(x) + Inches(0.22)), Emu(int(y) + Inches(0.72)), Emu(int(w) - Inches(0.35)), Inches(0.35),
             [(label, 12, lc, True)])
    if sub:
        add_text(slide, Emu(int(x) + Inches(0.22)), Emu(int(y) + Inches(1.05)), Emu(int(w) - Inches(0.35)), Inches(0.35),
                 [(sub, 10, sc, False)])


def bullet_block(slide, x, y, w, h, items, color=DARK, size=13, spacing=8):
    lines = []
    for t in items:
        lines.append((f"●  {t}", size, color, False, spacing, 1.15))
    return add_text(slide, x, y, w, h, lines)


def section_divider(no, en, title, subtitle, photo=None):
    s = add_slide(NAVY)
    add_rect(s, Inches(8.0), 0, Inches(5.4), SH, NAVY_2)
    if photo and Path(photo).exists():
        add_pic(s, photo, Inches(8.35), Inches(1.35), w=Inches(4.5), h=Inches(4.6))
        add_rect(s, Inches(8.35), Inches(1.35), Inches(4.5), Inches(0.08), RED)
        add_rect(s, Inches(8.35), Inches(5.87), Inches(4.5), Inches(0.08), GOLD)

    add_rect(s, Inches(0.7), Inches(2.0), Inches(0.12), Inches(1.6), RED)
    add_text(s, Inches(1.05), Inches(1.7), Inches(7), Inches(0.4),
             [(en, 12, GOLD, False)])
    add_text(s, Inches(1.0), Inches(2.15), Inches(7), Inches(1.0),
             [(f"{no}  {title}", 36, WHITE, True)])
    add_text(s, Inches(1.05), Inches(3.35), Inches(7), Inches(1.2),
             [(subtitle, 15, MUTED, False, 6, 1.25)])
    add_text(s, Inches(1.05), Inches(6.7), Inches(7), Inches(0.35),
             [("复旦科技园创业孵化基地  ·  领导汇报材料", 11, GREY, False)])
    return s


# ============================================================================
# 1. 封面
# ============================================================================
s = add_slide(WHITE)
# 右侧建筑大图
add_pic(s, PHOTOS / "building_innovation_center.jpeg", Inches(6.5), 0, h=SH)
# 左侧白底遮罩曲线效果：宽白区
add_rect(s, 0, 0, Inches(7.35), SH, WHITE)
# 斜切过渡
add_rect(s, Inches(6.4), 0, Inches(1.6), SH, WHITE, shape=MSO_SHAPE.RIGHT_TRIANGLE)
# 底部海军蓝条
add_rect(s, 0, Inches(6.85), SW, Inches(0.65), NAVY)
# 四圆点装饰（对标 2023 封面）
for i in range(4):
    add_rect(s, Inches(0.55 + i * 0.22), Inches(7.08), Inches(0.1), Inches(0.1), WHITE, shape=MSO_SHAPE.OVAL)
    add_rect(s, Inches(4.2 + i * 0.22), Inches(7.08), Inches(0.1), Inches(0.1), WHITE, shape=MSO_SHAPE.OVAL)

add_pic(s, ASSETS / "logo_stack.png", Inches(0.7), Inches(0.55), h=Inches(1.55))
add_text(s, Inches(2.45), Inches(0.75), Inches(4.5), Inches(1.1),
         [("复旦大学", 22, DARK, True, 2),
          ("国家大学科技园", 16, GREY, False)])

add_rect(s, Inches(0.75), Inches(2.55), Inches(1.1), Inches(0.06), RED)
add_text(s, Inches(0.7), Inches(2.85), Inches(6.2), Inches(1.6),
         [("2026 年度工作报告", 40, NAVY, True, 8),
          ("上海市创业孵化示范基地", 20, DARK, False)])
add_text(s, Inches(0.72), Inches(4.85), Inches(6), Inches(1.0),
         [("复旦科技园创业孵化基地", 16, STEEL, True, 4),
          ("聚焦成果转化 · 厚植创业生态 · 赋能区域发展", 13, GREY, False)])
add_text(s, Inches(0.72), Inches(6.2), Inches(5), Inches(0.35),
         [("汇报单位：复旦大学国家大学科技园    2026.01", 12, GREY, False)])
add_text(s, Inches(9.2), Inches(7.0), Inches(3.7), Inches(0.35),
         [("杨浦 · 国泰路 11 号", 11, MUTED, False)], align=PP_ALIGN.RIGHT)


# ============================================================================
# 2. 目录
# ============================================================================
s = add_slide(OFF)
add_rect(s, 0, 0, SW, Inches(0.9), NAVY)
add_pic(s, ASSETS / "logo_stack.png", Inches(0.35), Inches(0.12), h=Inches(0.66))
add_text(s, Inches(1.7), Inches(0.25), Inches(8), Inches(0.45),
         [("目  录  CONTENTS", 22, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, 0, Inches(0.9), SW, Inches(0.035), RED)

toc = [
    ("01", "载体基本情况", "平台定位 · 空间载体 · 荣誉资质 · 产业布局"),
    ("02", "2025 年工作成效", "孵化实效 · 就业带动 · 政策金融 · 服务活动"),
    ("03", "特色与亮点", "品牌建设 · 成果转化体系 · 区域协同开放"),
    ("04", "综合效益与展望", "经济 · 社会 · 绿色效益及下一步方向"),
]
colors = [RGBColor(0x3A, 0x7C, 0xB0), RGBColor(0x2A, 0x8F, 0x8A), RGBColor(0x2F, 0x5F, 0x8A), RGBColor(0x5A, 0x4E, 0x8A)]
for i, ((no, title, sub), col) in enumerate(zip(toc, colors)):
    y = Inches(1.45 + i * 1.25)
    add_rect(s, Inches(1.2), y, Inches(10.8), Inches(1.05), WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE).adjustments[0] = 0.1
    add_rect(s, Inches(1.2), y, Inches(0.18), Inches(1.05), col)
    diamond = add_rect(s, Inches(1.7), Emu(int(y) + Inches(0.22)), Inches(0.6), Inches(0.6), col, shape=MSO_SHAPE.DIAMOND)
    tf = diamond.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = no
    _set_font(r, 14, WHITE, True)
    add_text(s, Inches(2.7), Emu(int(y) + Inches(0.15)), Inches(8.5), Inches(0.4),
             [(title, 22, DARK, True)])
    add_text(s, Inches(2.7), Emu(int(y) + Inches(0.55)), Inches(8.5), Inches(0.35),
             [(sub, 13, GREY, False)])
footer(s, "02")


# ============================================================================
# 3. 开篇导语
# ============================================================================
s = add_slide(OFF)
header_bar(s, "导语", "总体定位与工作主线")
add_rect(s, Inches(0.55), Inches(1.15), Inches(12.2), Inches(1.55), NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE).adjustments[0] = 0.06
add_text(
    s,
    Inches(0.85),
    Inches(1.3),
    Inches(11.6),
    Inches(1.3),
    [
        ("紧扣国家大学科技园优化重塑主线，围绕", 15, WHITE, False, 2),
        ("“成果发现－概念验证－创业孵化－产业培育”", 18, GOLD, True, 4),
        ("推动高校成果加速转化、创业项目落地成长、科技企业集聚发展。", 15, WHITE, False),
    ],
)

pillars = [
    ("成果转化", "打通实验室到市场\n最初一公里"),
    ("创业孵化", "全链条培育\n科技型初创企业"),
    ("人才培养", "双创教育与\n实践赋能并进"),
    ("区域协同", "校地融合服务\n高质量发展"),
]
for i, (t, d) in enumerate(pillars):
    x = Inches(0.55 + i * 3.15)
    card = add_rect(s, x, Inches(3.05), Inches(2.95), Inches(2.7), WHITE, line_color=LINE, line_w=0.8, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    card.adjustments[0] = 0.08
    add_rect(s, x, Inches(3.05), Inches(2.95), Inches(0.1), RED)
    add_text(s, x, Inches(3.35), Inches(2.95), Inches(0.5), [(f"0{i+1}", 20, RED, True)], align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(3.9), Inches(2.95), Inches(0.5), [(t, 18, NAVY, True)], align=PP_ALIGN.CENTER)
    add_text(s, Emu(int(x) + Inches(0.2)), Inches(4.55), Inches(2.55), Inches(0.9),
             [(d, 13, GREY, False, 2, 1.25)], align=PP_ALIGN.CENTER)
footer(s, "03")


# ============================================================================
# 4. 章节分隔：载体基本情况
# ============================================================================
section_divider(
    "01",
    "PART ONE",
    "载体基本情况",
    "高校校区 · 科技基地 · 公共社区“三区联动”创新实践基地",
    PHOTOS / "building_innovation_center.jpeg",
)


# ============================================================================
# 5. 载体概况
# ============================================================================
s = add_slide(OFF)
header_bar(s, "01", "载体基本情况 · 平台定位与空间载体")

add_pic(s, PHOTOS / "incubator_interior.jpeg", Inches(0.5), Inches(1.15), w=Inches(5.1))
add_rect(s, Inches(0.5), Inches(5.55), Inches(5.1), Inches(1.35), NAVY)
add_text(
    s,
    Inches(0.7),
    Inches(5.7),
    Inches(4.7),
    Inches(1.1),
    [
        ("国泰路 11 号 · 场地面积 3.4 万㎡", 14, GOLD, True, 4),
        ("孵化面积 2.56 万㎡，占比 75%\n紧邻复旦邯郸校区与五角场商圈", 12, WHITE, False, 2, 1.2),
    ],
)

add_text(s, Inches(5.95), Inches(1.15), Inches(6.8), Inches(0.45),
         [("复旦大学国家大学科技园 · 首批国家大学科技园", 16, NAVY, True)])
add_text(
    s,
    Inches(5.95),
    Inches(1.65),
    Inches(6.8),
    Inches(1.6),
    [
        ("创立于 2000 年，由科技部、教育部联合认定。", 13, DARK, False, 6, 1.2),
        ("发展理念：“转化科技、服务社会、汇聚智慧、共创未来”。", 13, DARK, False, 6, 1.2),
        ("核心功能：科技成果转化、创业孵化、产业培育、创新人才培养。", 13, DARK, False, 6, 1.2),
    ],
)

facts = [
    ("全链条矩阵", "众创空间 + 孵化器\n+ 加速器 + 产业基地"),
    ("服务对象", "高校师生 / 校友团队\n科技型初创企业"),
    ("服务体系", "空间 · 政策 · 辅导\n技术转移 · 投融资"),
]
for i, (t, d) in enumerate(facts):
    y = Inches(3.55 + i * 1.05)
    add_rect(s, Inches(5.95), y, Inches(6.8), Inches(0.95), WHITE, line_color=LINE, line_w=0.75)
    add_rect(s, Inches(5.95), y, Inches(0.1), Inches(0.95), RED)
    add_text(s, Inches(6.3), Emu(int(y) + Inches(0.12)), Inches(2.2), Inches(0.7),
             [(t, 14, NAVY, True)], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(8.6), Emu(int(y) + Inches(0.12)), Inches(3.9), Inches(0.7),
             [(d, 12, GREY, False, 2, 1.15)], anchor=MSO_ANCHOR.MIDDLE)
footer(s, "05")


# ============================================================================
# 6. 荣誉与产业
# ============================================================================
s = add_slide(OFF)
header_bar(s, "01", "载体基本情况 · 荣誉资质与产业集聚")

honors = [
    "国家高新技术创业服务中心",
    "国家高校学生科技创业实习基地",
    "上海市创业孵化示范基地",
    "上海市海聚英才创新创业示范基地",
    "上海市知识产权示范基地",
    "上海市科技创新创业服务站",
    "中国高校孵化器十强",
    "国家火炬计划先进服务机构",
]
add_text(s, Inches(0.55), Inches(1.1), Inches(6), Inches(0.4), [("核心荣誉资质", 16, NAVY, True)])
for i, h in enumerate(honors):
    col = i % 2
    row = i // 2
    x = Inches(0.55 + col * 3.15)
    y = Inches(1.55 + row * 0.72)
    add_rect(s, x, y, Inches(3.0), Inches(0.6), WHITE, line_color=LINE, line_w=0.7, shape=MSO_SHAPE.ROUNDED_RECTANGLE).adjustments[0] = 0.15
    add_rect(s, x, y, Inches(0.08), Inches(0.6), RED)
    add_text(s, Emu(int(x) + Inches(0.2)), y, Inches(2.7), Inches(0.6),
             [(h, 11, DARK, False)], anchor=MSO_ANCHOR.MIDDLE)

# 右侧累计数据
add_rect(s, Inches(7.0), Inches(1.1), Inches(5.8), Inches(5.65), NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE).adjustments[0] = 0.04
add_text(s, Inches(7.3), Inches(1.3), Inches(5.2), Inches(0.4),
         [("截至 2025 年 · 关键培育数据", 15, GOLD, True)])
stats = [
    ("1500+", "存续注册企业", "科技型企业占比 80%"),
    ("115", "高新技术企业", "累计培育"),
    ("28 / 1", "专精特新中小企业 / 小巨人", "上海市级"),
    ("14", "资本市场挂牌上市企业", "多层次资本市场"),
    ("1.37 亿", "区级税收贡献", "2025 年基地企业"),
]
for i, (v, l, sub) in enumerate(stats):
    y = Inches(1.85 + i * 0.9)
    add_text(s, Inches(7.4), y, Inches(2.0), Inches(0.7), [(v, 22, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(9.4), y, Inches(3.1), Inches(0.7),
             [(l, 13, WHITE, True, 2), (sub, 11, MUTED, False)], anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.55), Inches(4.65), Inches(6.2), Inches(0.35),
         [("特色产业集群", 14, NAVY, True)])
clusters = ["集成电路", "软件和信息技术", "节能环保", "新材料新能源", "生物医药"]
for i, c in enumerate(clusters):
    x = Inches(0.55 + i * 1.25)
    chip = add_rect(s, x, Inches(5.1), Inches(1.15), Inches(0.48), NAVY_2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    chip.adjustments[0] = 0.3
    add_text(s, x, Inches(5.1), Inches(1.15), Inches(0.48),
             [(c, 10, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(0.55), Inches(5.8), Inches(6.2), Inches(0.9),
         [("初步形成以硬科技为主导、多领域协同发展的产业生态，\n为区域创新策源与产业升级提供持续动能。", 12, GREY, False, 2, 1.25)])
footer(s, "06")


# ============================================================================
# 7. 章节分隔：工作成效
# ============================================================================
section_divider(
    "02",
    "PART TWO",
    "2025 年工作成效",
    "以实效为导向，系统提升孵化服务能力与创业生态活跃度",
    PHOTOS / "meeting_alumni.jpeg",
)


# ============================================================================
# 8. 成效总览 KPI
# ============================================================================
s = add_slide(OFF)
header_bar(s, "02", "2025 年工作成效 · 关键指标总览")

kpis = [
    ("26 家", "新增入孵创业团队", "孵化成功率 100%"),
    ("159 家", "年末在孵创业实体", "初创组织存活率 100%"),
    ("324 人", "在孵组织在岗人数", "同比增长 57.28%"),
    ("9700 万+", "投融资支持金额", "助力企业加速成长"),
    ("47 场", "创新创业活动", "其中人社合作 3 场"),
    ("21 家", "新增高企（含复审）", "梯度培育持续见效"),
]
for i, (v, l, sub) in enumerate(kpis):
    col, row = i % 3, i // 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.2 + row * 2.55)
    card = add_rect(s, x, y, Inches(4.0), Inches(2.3), WHITE, line_color=LINE, line_w=0.8, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    card.adjustments[0] = 0.07
    add_rect(s, x, y, Inches(4.0), Inches(0.1), RED if row == 0 else NAVY)
    add_text(s, x, Emu(int(y) + Inches(0.45)), Inches(4.0), Inches(0.7),
             [(v, 32, RED if i % 2 == 0 else NAVY, True)], align=PP_ALIGN.CENTER)
    add_text(s, x, Emu(int(y) + Inches(1.2)), Inches(4.0), Inches(0.4),
             [(l, 15, DARK, True)], align=PP_ALIGN.CENTER)
    add_text(s, x, Emu(int(y) + Inches(1.65)), Inches(4.0), Inches(0.35),
             [(sub, 12, GREY, False)], align=PP_ALIGN.CENTER)
footer(s, "08")


# ============================================================================
# 9. 创业孵化实效
# ============================================================================
s = add_slide(OFF)
header_bar(s, "02", "创业孵化实效 · 项目挖掘与企业培育")

# 左栏：项目挖掘
add_rect(s, Inches(0.45), Inches(1.1), Inches(6.15), Inches(5.65), WHITE, line_color=LINE, line_w=0.75)
add_rect(s, Inches(0.45), Inches(1.1), Inches(6.15), Inches(0.55), NAVY)
add_text(s, Inches(0.65), Inches(1.18), Inches(5.7), Inches(0.4),
         [("一、加强创业项目挖掘", 15, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.7), Inches(1.85), Inches(5.7), Inches(0.7),
         [("聚焦人工智能、集成电路、生物医药、新材料等重点领域，\n与院系、科研团队和创新创业平台协同挖掘产业化项目。", 12, GREY, False, 2, 1.2)])

mini = [
    ("16", "对接院系（个）"),
    ("21", "新增入库成果"),
    ("12", "高校成果转化项"),
    ("11", "深度服务衍生企业"),
]
for i, (v, l) in enumerate(mini):
    x = Inches(0.7 + (i % 2) * 2.9)
    y = Inches(2.75 + (i // 2) * 1.15)
    add_rect(s, x, y, Inches(2.7), Inches(1.0), OFF)
    add_text(s, x, Emu(int(y) + Inches(0.1)), Inches(2.7), Inches(0.45),
             [(v, 24, RED, True)], align=PP_ALIGN.CENTER)
    add_text(s, x, Emu(int(y) + Inches(0.55)), Inches(2.7), Inches(0.35),
             [(l, 11, DARK, False)], align=PP_ALIGN.CENTER)

add_text(s, Inches(0.7), Inches(5.2), Inches(5.7), Inches(1.2),
         [("重点服务：益臻新能源、波达医疗、眸深智能、复泓智研等\n复旦科创衍生企业，提供知识产权、政策、融资、落地等定制支持。",
           12, DARK, False, 4, 1.25)])

# 右栏：企业培育
add_rect(s, Inches(6.85), Inches(1.1), Inches(6.0), Inches(5.65), WHITE, line_color=LINE, line_w=0.75)
add_rect(s, Inches(6.85), Inches(1.1), Inches(6.0), Inches(0.55), NAVY_2)
add_text(s, Inches(7.05), Inches(1.18), Inches(5.6), Inches(0.4),
         [("二、提升企业培育能力", 15, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)

bullet_block(
    s,
    Inches(7.1),
    Inches(1.9),
    Inches(5.5),
    Inches(2.4),
    [
        "依托“众创空间+孵化器+加速器+产业基地”矩阵综合赋能",
        "新增入孵创业团队 26 家，孵化成功率 100%",
        "新增入孵初创组织 23 家，年末存活率 100%",
        "年末在孵创业实体 159 家",
        "建立科技型中小企业—高企—领军企业梯度培育机制",
    ],
    size=12,
    spacing=7,
)

grad = [
    ("21", "高新技术企业\n（含复审）"),
    ("2", "杨浦科技小巨人\n（含培育）"),
    ("3", "“3310”企业"),
]
for i, (v, l) in enumerate(grad):
    x = Inches(7.15 + i * 1.85)
    add_rect(s, x, Inches(4.7), Inches(1.7), Inches(1.6), NAVY)
    add_text(s, x, Inches(4.85), Inches(1.7), Inches(0.55),
             [(v, 26, GOLD, True)], align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(5.5), Inches(1.7), Inches(0.65),
             [(l, 11, WHITE, False, 1, 1.15)], align=PP_ALIGN.CENTER)
footer(s, "09")


# ============================================================================
# 10. 典型案例选树
# ============================================================================
s = add_slide(OFF)
header_bar(s, "02", "创业孵化实效 · 培育创业典型")

add_text(s, Inches(0.55), Inches(1.1), Inches(12), Inches(0.4),
         [("2025 年产生创业评选优秀创业组织 9 家，形成一批具有高校特色的转化与创业标杆。", 13, GREY, False)])

cases = [
    ("眸深智能", "全国颠覆性技术创新大赛\n总决赛最高奖及多项赛事奖项", RED),
    ("波达医疗", "国家级/市级科创赛事优异成绩\n获批市关键技术研发计划项目", NAVY),
    ("赤子青山团队", "中国国际大学生创新大赛\n“青年筑梦红色之旅”金奖", NAVY_2),
    ("音书科技", "2025 年度\n“天使基金优秀雏鹰企业”", STEEL),
]
for i, (name, desc, col) in enumerate(cases):
    x = Inches(0.5 + i * 3.2)
    card = add_rect(s, x, Inches(1.7), Inches(3.0), Inches(3.35), WHITE, line_color=LINE, line_w=0.8, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    card.adjustments[0] = 0.07
    add_rect(s, x, Inches(1.7), Inches(3.0), Inches(0.9), col)
    add_text(s, x, Inches(1.85), Inches(3.0), Inches(0.6),
             [(name, 18, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Emu(int(x) + Inches(0.2)), Inches(2.9), Inches(2.6), Inches(1.8),
             [(desc, 13, DARK, False, 4, 1.3)], align=PP_ALIGN.CENTER)

add_rect(s, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.4), NAVY)
add_text(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.05),
         [("同步成效：3 个创业团队获上海市大学生科技创业基金复旦大学基金支持；", 13, WHITE, False, 4),
          ("波达医疗联合复旦大学团队获批上海市 2025 年度关键技术研发计划项目，体现“师生共创、园区赋能”路径。", 13, GOLD, False)],
         )
footer(s, "10")


# ============================================================================
# 11. 就业 + 政策金融
# ============================================================================
s = add_slide(OFF)
header_bar(s, "02", "创业带动就业 · 政策落实与金融支持")

# 就业
add_rect(s, Inches(0.45), Inches(1.1), Inches(4.1), Inches(5.65), WHITE, line_color=LINE, line_w=0.75)
add_rect(s, Inches(0.45), Inches(1.1), Inches(4.1), Inches(0.55), RED)
add_text(s, Inches(0.65), Inches(1.18), Inches(3.7), Inches(0.4),
         [("创业带动就业", 15, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(0.7), Inches(1.9), Inches(3.6), Inches(0.9),
         [("低成本、便利化、专业化服务支撑企业扩岗；\n分层分类开展创新创业教育与实训。", 12, GREY, False, 2, 1.2)])
add_text(s, Inches(0.7), Inches(2.95), Inches(3.6), Inches(0.7),
         [("324", 40, RED, True)], align=PP_ALIGN.CENTER)
add_text(s, Inches(0.7), Inches(3.65), Inches(3.6), Inches(0.4),
         [("在孵创业组织在岗人数", 13, DARK, True)], align=PP_ALIGN.CENTER)
add_text(s, Inches(0.7), Inches(4.15), Inches(3.6), Inches(0.35),
         [("同比增长 57.28%", 14, NAVY, True)], align=PP_ALIGN.CENTER)
add_rect(s, Inches(0.85), Inches(4.75), Inches(3.3), Inches(1.5), OFF)
add_text(s, Inches(1.0), Inches(4.95), Inches(3.0), Inches(1.15),
         [("新增就业 118 人", 16, NAVY, True, 6),
          ("为高校毕业生、青年创业者与专业技术人才提供岗位机会", 11, GREY, False, 2, 1.2)])

# 政策
add_rect(s, Inches(4.75), Inches(1.1), Inches(4.1), Inches(5.65), WHITE, line_color=LINE, line_w=0.75)
add_rect(s, Inches(4.75), Inches(1.1), Inches(4.1), Inches(0.55), NAVY)
add_text(s, Inches(4.95), Inches(1.18), Inches(3.7), Inches(0.4),
         [("创业政策落实", 15, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
bullet_block(
    s,
    Inches(5.0),
    Inches(1.9),
    Inches(3.6),
    Inches(3.2),
    [
        "高企认定 / 科技项目 / “3310” / 人社扶持辅导",
        "新增高企（含复审）21 家",
        "杨浦科技小巨人（含培育）2 家",
        "“3310”企业 3 家",
        "推荐 9 家享受人社创业扶持政策",
        "帮扶 2 家获得创业担保贷款",
    ],
    size=12,
    spacing=6,
)
add_pic(s, PHOTOS / "site_visit.jpeg", Inches(5.0), Inches(5.15), w=Inches(3.6))

# 金融
add_rect(s, Inches(9.05), Inches(1.1), Inches(3.8), Inches(5.65), NAVY)
add_text(s, Inches(9.25), Inches(1.25), Inches(3.4), Inches(0.4),
         [("多层次金融支持", 15, GOLD, True)])
add_text(s, Inches(9.25), Inches(1.8), Inches(3.4), Inches(1.5),
         [("投资机构合作", 12, MUTED, False, 4),
          ("山东工研基金、禹泽资本、小村创投、创业接力基金、杨浦科创集团等", 12, WHITE, False, 4, 1.2)])
add_text(s, Inches(9.25), Inches(3.4), Inches(3.4), Inches(1.3),
         [("银行机构合作", 12, MUTED, False, 4),
          ("上海农商行、招商银行、中国银行等融资服务合作", 12, WHITE, False, 4, 1.2)])
add_rect(s, Inches(9.25), Inches(4.85), Inches(3.4), Inches(1.5), NAVY_2)
add_text(s, Inches(9.4), Inches(5.0), Inches(3.1), Inches(1.25),
         [("帮扶获投融资约", 12, MUTED, False, 2),
          ("9700 余万元", 24, GOLD, True, 2),
          ("联动大学生创业基金与科创母基金", 11, WHITE, False)])
footer(s, "11")


# ============================================================================
# 12. 孵化服务体系
# ============================================================================
s = add_slide(OFF)
header_bar(s, "02", "创业孵化服务 · 空间、专业服务与导师体系")

cols = [
    ("空间支持", NAVY, [
        "综合更新与功能提升项目获批张江专项重点项目",
        "场地 3.41 万㎡，孵化面积占比 75%",
        "向 49 家创业组织提供租金减免等支持",
        "办公 / 共享 / 会议路演空间一体化配置",
    ]),
    ("专业服务", RED, [
        "政策、知识产权、法律财税、投融资、市场拓展",
        "签约入驻第三方服务机构 41 个",
        "全年提供服务超过 2000 家次",
        "落实创业帮扶指标 40 家",
    ]),
    ("导师赋能", STEEL, [
        "新增创业导师 16 人",
        "高校教师 + 企业管理者 + 投资专家",
        "覆盖技术、管理、资本、市场辅导",
        "一对一辅导与主题沙龙并重",
    ]),
]
for i, (title, col, items) in enumerate(cols):
    x = Inches(0.45 + i * 4.25)
    add_rect(s, x, Inches(1.15), Inches(4.05), Inches(5.55), WHITE, line_color=LINE, line_w=0.75)
    add_rect(s, x, Inches(1.15), Inches(4.05), Inches(0.7), col)
    add_text(s, x, Inches(1.25), Inches(4.05), Inches(0.5),
             [(title, 18, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    bullet_block(s, Emu(int(x) + Inches(0.25)), Inches(2.15), Inches(3.55), Inches(4.2), items, size=13, spacing=10)
footer(s, "12")


# ============================================================================
# 13. 复煜概念验证中心
# ============================================================================
s = add_slide(OFF)
header_bar(s, "02", "服务亮点 · 复煜概念验证中心正式运营")

add_rect(s, Inches(0.45), Inches(1.15), Inches(7.6), Inches(5.55), WHITE, line_color=LINE, line_w=0.75)
add_text(s, Inches(0.75), Inches(1.4), Inches(7.0), Inches(0.45),
         [("打通科技成果转化“最初一公里”", 20, NAVY, True)])
add_text(
    s,
    Inches(0.75),
    Inches(2.0),
    Inches(7.0),
    Inches(1.2),
    [
        ("复煜概念验证中心于 2025 年 8 月投入运营，围绕新材料、人工智能、集成电路等领域开展首批项目征集，", 13, DARK, False, 4, 1.25),
        ("从 60 余个项目中筛选 4 个项目进入验证阶段，提供资金、载体、应用场景、科研支撑、投融资对接与运营协助。", 13, DARK, False, 4, 1.25),
    ],
)

flow = ["项目征集", "专业遴选", "概念验证", "孵化加速"]
for i, t in enumerate(flow):
    x = Inches(0.85 + i * 1.8)
    chip = add_rect(s, x, Inches(3.5), Inches(1.5), Inches(0.7), NAVY if i < 3 else RED, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    chip.adjustments[0] = 0.25
    add_text(s, x, Inches(3.5), Inches(1.5), Inches(0.7),
             [(t, 13, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < 3:
        add_text(s, Emu(int(x) + Inches(1.45)), Inches(3.55), Inches(0.4), Inches(0.6),
                 [("→", 18, GOLD, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

metrics = [("60+", "征集项目"), ("4", "进入验证"), ("6 维", "全方位服务")]
for i, (v, l) in enumerate(metrics):
    x = Inches(0.85 + i * 2.35)
    add_rect(s, x, Inches(4.6), Inches(2.15), Inches(1.6), OFF)
    add_text(s, x, Inches(4.75), Inches(2.15), Inches(0.7),
             [(v, 26, RED, True)], align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(5.5), Inches(2.15), Inches(0.5),
             [(l, 13, DARK, False)], align=PP_ALIGN.CENTER)

add_pic(s, PHOTOS / "tech_transfer.jpeg", Inches(8.3), Inches(1.15), w=Inches(4.5))
add_rect(s, Inches(8.3), Inches(4.55), Inches(4.5), Inches(2.15), NAVY)
add_text(s, Inches(8.55), Inches(4.75), Inches(4.1), Inches(1.8),
         [("服务内涵", 14, GOLD, True, 6),
          ("技术验证 · 应用场景 · 科研支撑\n载体空间 · 投融资对接 · 运营辅导", 13, WHITE, False, 4, 1.3)])
footer(s, "13")


# ============================================================================
# 14. 活动组织
# ============================================================================
s = add_slide(OFF)
header_bar(s, "02", "创新创业活动组织 · 课程、产学研与赛事")

# 三块
blocks = [
    ("完善双创课程体系", [
        "“复·创课堂”“复旦科技园企业沙龙”聚焦创业实务",
        "马兰花计划 GYB 培训班顺利举办",
        "联合创新创业学院开发“理论+实践”课程",
    ]),
    ("强化产学研交流", [
        "人工智能产业应用路演",
        "新材料研发与产业应用研讨会",
        "无掩膜光刻微纳器件加工产学研沙龙",
    ]),
    ("以赛聚才引才", [
        "复旦科创大赛 / 复旦之星",
        "创·在上海 / 杨浦“科创之星”",
        "推荐参赛项目 100 余项",
    ]),
]
for i, (title, items) in enumerate(blocks):
    x = Inches(0.45 + i * 4.25)
    add_rect(s, x, Inches(1.1), Inches(4.05), Inches(2.85), WHITE, line_color=LINE, line_w=0.75)
    add_rect(s, x, Inches(1.1), Inches(4.05), Inches(0.5), NAVY if i != 1 else RED)
    add_text(s, x, Inches(1.15), Inches(4.05), Inches(0.4),
             [(title, 14, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    bullet_block(s, Emu(int(x) + Inches(0.2)), Inches(1.8), Inches(3.65), Inches(2.0), items, size=12, spacing=7)

# 照片条
photos_row = [
    PHOTOS / "contest1.jpeg",
    PHOTOS / "activity_salon.jpeg",
    PHOTOS / "edu_activity1.jpeg",
    PHOTOS / "practice1.jpeg",
]
for i, p in enumerate(photos_row):
    add_pic(s, p, Inches(0.45 + i * 3.2), Inches(4.2), w=Inches(3.05), h=Inches(1.7))

add_rect(s, Inches(0.45), Inches(6.05), Inches(12.4), Inches(0.75), NAVY)
add_text(s, Inches(0.7), Inches(6.15), Inches(12), Inches(0.55),
         [("全年举办创新创业活动 47 场（含人社合作 3 场）；推荐人社类活动/赛事 12 项，持续提升生态活跃度。", 13, WHITE, False)],
         anchor=MSO_ANCHOR.MIDDLE)
footer(s, "14")


# ============================================================================
# 15. 复旦科创大赛专题
# ============================================================================
s = add_slide(OFF)
header_bar(s, "02", "品牌活动专题 · 复旦科创大赛")

add_pic(s, PHOTOS / "contest2.jpeg", Inches(0.45), Inches(1.15), w=Inches(5.4), h=Inches(3.2))
add_pic(s, PHOTOS / "contest3.jpeg", Inches(0.45), Inches(4.5), w=Inches(5.4), h=Inches(2.2))

add_text(s, Inches(6.15), Inches(1.15), Inches(6.6), Inches(0.45),
         [("从展示平台迈向组织平台", 18, NAVY, True)])
add_text(
    s,
    Inches(6.15),
    Inches(1.7),
    Inches(6.6),
    Inches(1.3),
    [
        ("由复旦大学、虹口区人民政府、锦江国际集团共同指导，基地作为主要承办单位之一，承担赛事组织、项目遴选、创业辅导、评审协调、资源对接及赛后孵化。", 12, DARK, False, 4, 1.25),
    ],
)

race_kpis = [
    ("456", "报名项目"),
    ("20", "晋级决赛团队"),
    ("双赛道", "创意组 + 创业组"),
]
for i, (v, l) in enumerate(race_kpis):
    x = Inches(6.15 + i * 2.2)
    add_rect(s, x, Inches(3.2), Inches(2.05), Inches(1.35), WHITE, line_color=LINE, line_w=0.7)
    add_text(s, x, Inches(3.35), Inches(2.05), Inches(0.6),
             [(v, 22, RED, True)], align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(3.95), Inches(2.05), Inches(0.4),
             [(l, 12, DARK, False)], align=PP_ALIGN.CENTER)

add_rect(s, Inches(6.15), Inches(4.8), Inches(6.6), Inches(1.9), NAVY)
add_text(s, Inches(6.4), Inches(5.0), Inches(6.2), Inches(1.55),
         [("评审维度：技术原创性 － 工程可行性 － 转化现实性", 13, GOLD, True, 6),
          ("覆盖信息技术、集成电路、生命健康、新材料等重点领域；部分项目已纳入成果转化资源库并获配套支持，推动创新链与产业链深度对接。", 12, WHITE, False, 2, 1.25)])
footer(s, "15")


# ============================================================================
# 16. 媒体宣传
# ============================================================================
s = add_slide(OFF)
header_bar(s, "02", "创业媒体宣传 · 品牌焕新与典型传播")

add_rect(s, Inches(0.45), Inches(1.15), Inches(8.0), Inches(5.55), WHITE, line_color=LINE, line_w=0.75)
add_text(s, Inches(0.75), Inches(1.4), Inches(7.5), Inches(0.4),
         [("参与复旦科技园品牌焕新工程", 18, NAVY, True)])
bullet_block(
    s,
    Inches(0.8),
    Inches(2.0),
    Inches(7.3),
    Inches(2.4),
    [
        "协助完成园区官网、微信公众号等宣传阵地升级迭代",
        "挖掘品牌内涵、提升信息质量，塑造契合复旦特色的科创形象",
        "重点宣传眸深智能、波达医疗等成长案例与赛事优秀项目",
        "增强高校、产业、资本与区域创新主体之间的资源链接能力",
    ],
    size=13,
    spacing=8,
)

pub = [
    ("34", "官网/公号发布稿件"),
    ("11", "向人社部门供稿"),
    ("10", "被海纳百创等采纳"),
]
for i, (v, l) in enumerate(pub):
    x = Inches(0.8 + i * 2.5)
    add_rect(s, x, Inches(4.7), Inches(2.3), Inches(1.5), OFF)
    add_text(s, x, Inches(4.85), Inches(2.3), Inches(0.65),
             [(v, 28, RED, True)], align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(5.55), Inches(2.3), Inches(0.45),
             [(l, 12, DARK, False)], align=PP_ALIGN.CENTER)

add_pic(s, PHOTOS / "activity_lecture.jpeg", Inches(8.7), Inches(1.15), w=Inches(4.1), h=Inches(2.6))
add_pic(s, PHOTOS / "workshop.jpeg", Inches(8.7), Inches(3.95), w=Inches(4.1), h=Inches(2.75))
footer(s, "16")


# ============================================================================
# 17. 章节分隔：特色亮点
# ============================================================================
section_divider(
    "03",
    "PART THREE",
    "特色与亮点",
    "成果转化型创业孵化模式：品牌 · 体系 · 生态",
    PHOTOS / "lab_visit.jpeg",
)


# ============================================================================
# 18. 特色一：品牌
# ============================================================================
s = add_slide(OFF)
header_bar(s, "03", "特色亮点（一）打造特色品牌，提升孵化影响力")

add_text(
    s,
    Inches(0.55),
    Inches(1.15),
    Inches(12.2),
    Inches(0.8),
    [
        ("以二十余年国家大学科技园建设积淀为根基，持续强化高校特色创业孵化品牌，", 14, DARK, False, 2),
        ("构建覆盖创业辅导、赛事、技术转移、政策、知产、财税、投融资、人才、市场的一站式服务网络。", 14, DARK, False),
    ],
)

brand_cards = [
    ("41 家", "第三方专业服务机构", "全年服务超 2000 家次"),
    ("16 人", "新增签约创业导师", "复合型导师团队扩容"),
    ("47 场", "创新创业活动", "持续激活园区生态"),
    ("190 家", "新增注册企业", "科技型企业 143 家"),
]
for i, (v, l, sub) in enumerate(brand_cards):
    x = Inches(0.5 + i * 3.2)
    add_rect(s, x, Inches(2.2), Inches(3.0), Inches(2.0), WHITE, line_color=LINE, line_w=0.75, shape=MSO_SHAPE.ROUNDED_RECTANGLE).adjustments[0] = 0.08
    add_text(s, x, Inches(2.4), Inches(3.0), Inches(0.7),
             [(v, 28, RED, True)], align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(3.15), Inches(3.0), Inches(0.4),
             [(l, 13, DARK, True)], align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(3.55), Inches(3.0), Inches(0.4),
             [(sub, 11, GREY, False)], align=PP_ALIGN.CENTER)

add_rect(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.2), NAVY)
add_text(s, Inches(0.8), Inches(4.7), Inches(7.5), Inches(1.8),
         [("区域贡献同步增强", 15, GOLD, True, 6),
          ("2025 年基地企业实现区级税收贡献 1.37 亿元；新增高新技术企业（含复审）21 家、杨浦科技小巨人（含培育）2 家、“3310”企业 3 家，品牌影响力与发展贡献力同向提升。", 13, WHITE, False, 4, 1.3)])
add_pic(s, PHOTOS / "recruitment_fair.jpeg", Inches(8.6), Inches(4.7), w=Inches(3.9), h=Inches(1.8))
footer(s, "18")


# ============================================================================
# 19. 特色二：服务模式 + 波达案例
# ============================================================================
s = add_slide(OFF)
header_bar(s, "03", "特色亮点（二）成果转化型孵化体系 · 波达医疗案例")

# 链条
chain = ["成果发现", "概念验证", "创业孵化", "产业培育"]
for i, t in enumerate(chain):
    x = Inches(0.5 + i * 3.2)
    add_rect(s, x, Inches(1.15), Inches(2.95), Inches(0.7), NAVY if i % 2 == 0 else NAVY_2, shape=MSO_SHAPE.ROUNDED_RECTANGLE).adjustments[0] = 0.2
    add_text(s, x, Inches(1.15), Inches(2.95), Inches(0.7),
             [(f"{i+1}. {t}", 14, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_rect(s, Inches(0.45), Inches(2.1), Inches(6.2), Inches(4.6), WHITE, line_color=LINE, line_w=0.75)
add_text(s, Inches(0.7), Inches(2.25), Inches(5.7), Inches(0.4),
         [("体系创新要点", 15, NAVY, True)])
bullet_block(
    s,
    Inches(0.75),
    Inches(2.75),
    Inches(5.6),
    Inches(3.6),
    [
        "成果转化前移：对接 16 个院系，入库成果 21 项",
        "复煜概念验证中心承接早期硬科技项目",
        "全链条载体协同：众创—孵化—加速—产业基地",
        "服务升级：由单一孵化转向转化、成长、产业培育协同",
        "金融接力：财政引导 + 社会资本 + 产业落地拨投结合",
    ],
    size=13,
    spacing=8,
)

add_rect(s, Inches(6.9), Inches(2.1), Inches(5.95), Inches(4.6), NAVY)
add_text(s, Inches(7.15), Inches(2.25), Inches(5.5), Inches(0.4),
         [("孵化案例：波达医疗", 16, GOLD, True)])
add_text(
    s,
    Inches(7.15),
    Inches(2.75),
    Inches(5.5),
    Inches(3.7),
    [
        ("复旦信息学院教师团队创办，2022 年入驻基地，聚焦超快超声成像产业化。", 12, WHITE, False, 6, 1.2),
        ("服务举措：创业辅导、3310 申报、载体落地湾谷、租金/社保支持、对接投资资源。", 12, WHITE, False, 6, 1.2),
        ("成长结果：海聚英才大赛二等奖；获飞图创投数千万元融资，估值突破亿元。", 12, GOLD, True, 6, 1.2),
        ("产业意义：有望突破高端医学超声国际技术壁垒，服务脑科学与医学超声科研。", 12, WHITE, False, 4, 1.2),
    ],
)
footer(s, "19")


# ============================================================================
# 20. 特色三：区域协同与国际化
# ============================================================================
s = add_slide(OFF)
header_bar(s, "03", "特色亮点（三）深化区域协同，构建开放融合生态")

# 左：校地协同
add_rect(s, Inches(0.45), Inches(1.1), Inches(6.15), Inches(3.4), WHITE, line_color=LINE, line_w=0.75)
add_rect(s, Inches(0.45), Inches(1.1), Inches(6.15), Inches(0.5), NAVY)
add_text(s, Inches(0.65), Inches(1.18), Inches(5.7), Inches(0.35),
         [("校地融合 · 服务杨浦创新体系建设", 14, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
bullet_block(
    s,
    Inches(0.7),
    Inches(1.8),
    Inches(5.7),
    Inches(2.4),
    [
        "推动学校成果、人才与区域产业需求精准对接",
        "探索“财政引导资金+社会资本接续+产业协同落地”",
        "助力益臻新能源对接杨浦成果转化支持政策",
        "2025 年基地企业区级税收贡献 1.37 亿元",
    ],
    size=12,
    spacing=7,
)

# 右：产学研网络
add_rect(s, Inches(6.8), Inches(1.1), Inches(6.05), Inches(3.4), WHITE, line_color=LINE, line_w=0.75)
add_rect(s, Inches(6.8), Inches(1.1), Inches(6.05), Inches(0.5), RED)
add_text(s, Inches(7.0), Inches(1.18), Inches(5.7), Inches(0.35),
         [("产学研合作网络持续拓展", 14, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
partners = ["祖泉研究院", "复旦技术转移中心", "同济技术转移中心", "张江磁谷", "上海国际化工创新中心", "迈科技", "国科新研"]
for i, p in enumerate(partners):
    x = Inches(7.05 + (i % 2) * 2.8)
    y = Inches(1.85 + (i // 2) * 0.55)
    add_rect(s, x, y, Inches(2.6), Inches(0.45), OFF, shape=MSO_SHAPE.ROUNDED_RECTANGLE).adjustments[0] = 0.2
    add_text(s, x, y, Inches(2.6), Inches(0.45),
             [(p, 11, DARK, False)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 底部国际化
add_rect(s, Inches(0.45), Inches(4.7), Inches(12.4), Inches(2.0), NAVY)
add_text(s, Inches(0.75), Inches(4.9), Inches(6.5), Inches(1.6),
         [("国际化交流与合作", 14, GOLD, True, 6),
          ("参与发起“北欧创新国际会客厅”；举办丝路华章·大使领创计划出海沙龙、“复创·大咖说”等活动；接待巴西、俄罗斯等国际参访团，提升开放合作能级与国际影响力。", 12, WHITE, False, 2, 1.25)])
add_pic(s, PHOTOS / "intl_exchange.jpeg", Inches(7.6), Inches(4.9), w=Inches(2.4), h=Inches(1.6))
add_pic(s, PHOTOS / "intl_delegation.jpeg", Inches(10.15), Inches(4.9), w=Inches(2.4), h=Inches(1.6))
footer(s, "20")


# ============================================================================
# 21. 章节分隔：综合效益
# ============================================================================
section_divider(
    "04",
    "PART FOUR",
    "综合效益与展望",
    "经济发展 · 社会服务 · 绿色运营，持续释放孵化综合效益",
    PHOTOS / "building_street.jpeg",
)


# ============================================================================
# 22. 综合效益三栏
# ============================================================================
s = add_slide(OFF)
header_bar(s, "04", "综合效益 · 经济、社会与绿色发展")

benefits = [
    ("经济效益", "产业集聚效应增强", [
        "新增注册企业 190 家（科技型 143 家）",
        "新增高企 21 家、小巨人 2 家、3310 企业 3 家",
        "新增知识产权 318 项（发明专利 42 项）",
        "区级税收贡献 1.37 亿元",
    ], NAVY),
    ("社会效益", "双创生态持续完善", [
        "推动高校成果转化 12 项",
        "新增创业导师 16 人",
        "培训 11 场、活动 37 场、推荐赛事 100+ 项",
        "市载体绩效：复旦科技园优良，复翼互联优秀",
    ], RED),
    ("绿色效益", "载体运营品质优化", [
        "景观照明、公共空间、设施设备更新提升",
        "综合更新项目获批张江专项重点项目",
        "大厦通过“无废城市细胞”建设评估",
        "获评区级“无废楼宇”",
    ], STEEL),
]
for i, (title, sub, items, col) in enumerate(benefits):
    x = Inches(0.4 + i * 4.3)
    add_rect(s, x, Inches(1.15), Inches(4.1), Inches(5.55), WHITE, line_color=LINE, line_w=0.75)
    add_rect(s, x, Inches(1.15), Inches(4.1), Inches(1.15), col)
    add_text(s, x, Inches(1.25), Inches(4.1), Inches(0.5),
             [(title, 20, WHITE, True)], align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(1.75), Inches(4.1), Inches(0.4),
             [(sub, 12, GOLD if col != RED else WHITE, False)], align=PP_ALIGN.CENTER)
    bullet_block(s, Emu(int(x) + Inches(0.25)), Inches(2.55), Inches(3.6), Inches(3.8), items, size=13, spacing=10)
footer(s, "22")


# ============================================================================
# 23. 考评与展望
# ============================================================================
s = add_slide(OFF)
header_bar(s, "04", "阶段评价与下一步工作方向")

add_text(s, Inches(0.55), Inches(1.15), Inches(6), Inches(0.4),
         [("阶段性评价结果", 16, NAVY, True)])
evals = [
    ("上海市创新创业载体绩效评价", "复旦科技园：优良\n复翼互联众创空间：优秀"),
    ("杨浦区科技园区考评", "复旦科技园：优秀"),
]
for i, (t, d) in enumerate(evals):
    x = Inches(0.5 + i * 4.0)
    add_rect(s, x, Inches(1.65), Inches(3.8), Inches(1.9), WHITE, line_color=LINE, line_w=0.75)
    add_rect(s, x, Inches(1.65), Inches(0.1), Inches(1.9), RED)
    add_text(s, Emu(int(x) + Inches(0.25)), Inches(1.8), Inches(3.3), Inches(0.5),
             [(t, 13, NAVY, True)])
    add_text(s, Emu(int(x) + Inches(0.25)), Inches(2.4), Inches(3.3), Inches(0.9),
             [(d, 14, DARK, False, 2, 1.25)])

add_rect(s, Inches(8.5), Inches(1.65), Inches(4.3), Inches(1.9), NAVY)
add_text(s, Inches(8.75), Inches(1.9), Inches(3.9), Inches(1.4),
         [("核心判断", 13, GOLD, True, 6),
          ("已形成“成果持续转化、企业梯度成长、产业能级提升”的发展格局。", 13, WHITE, False, 2, 1.25)])

add_text(s, Inches(0.55), Inches(3.85), Inches(12), Inches(0.4),
         [("下一步工作方向", 16, NAVY, True)])
next_steps = [
    ("深化转化前移", "强化院系协同与概念验证，提升早期硬科技承接能力"),
    ("做精孵化服务", "完善导师、金融、场景与梯度培育机制"),
    ("放大品牌赛事", "推动复旦科创大赛从展示平台迈向组织平台"),
    ("拓展开放协同", "加深校地校企合作与国际化资源链接"),
]
for i, (t, d) in enumerate(next_steps):
    x = Inches(0.5 + (i % 4) * 3.2)
    y = Inches(4.4)
    add_rect(s, x, y, Inches(3.05), Inches(2.2), WHITE, line_color=LINE, line_w=0.75)
    add_rect(s, x, y, Inches(3.05), Inches(0.08), RED)
    add_text(s, x, Emu(int(y) + Inches(0.25)), Inches(3.05), Inches(0.45),
             [(f"0{i+1}  {t}", 14, NAVY, True)], align=PP_ALIGN.CENTER)
    add_text(s, Emu(int(x) + Inches(0.15)), Emu(int(y) + Inches(0.85)), Inches(2.75), Inches(1.1),
             [(d, 12, GREY, False, 2, 1.25)], align=PP_ALIGN.CENTER)
footer(s, "23")


# ============================================================================
# 24. 封底
# ============================================================================
s = add_slide(WHITE)
add_pic(s, PHOTOS / "building_street.jpeg", Inches(5.8), 0, h=SH)
add_rect(s, 0, 0, Inches(6.8), SH, WHITE)
add_rect(s, Inches(5.9), 0, Inches(1.5), SH, WHITE, shape=MSO_SHAPE.RIGHT_TRIANGLE)
add_rect(s, 0, Inches(6.85), SW, Inches(0.65), NAVY)
for i in range(4):
    add_rect(s, Inches(0.55 + i * 0.22), Inches(7.08), Inches(0.1), Inches(0.1), WHITE, shape=MSO_SHAPE.OVAL)

add_pic(s, ASSETS / "logo_stack.png", Inches(0.75), Inches(0.7), h=Inches(1.5))
add_rect(s, Inches(0.8), Inches(2.55), Inches(1.0), Inches(0.06), RED)
add_text(s, Inches(0.75), Inches(2.9), Inches(5.5), Inches(1.8),
         [("育高校创业新苗，", 28, NAVY, True, 6),
          ("助技术成果转化！", 28, NAVY, True)])
add_text(s, Inches(0.75), Inches(5.0), Inches(5.5), Inches(1.0),
         [("复旦大学国家大学科技园", 16, DARK, True, 4),
          ("创业孵化基地 · 2026 年度工作报告", 13, GREY, False)])
add_text(s, Inches(0.75), Inches(6.35), Inches(5.5), Inches(0.3),
         [("感谢各位领导指导！", 14, STEEL, True)])


# ----------------------------------------------------------------------------
# 输出
# ----------------------------------------------------------------------------
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"已生成：{OUT}")
print(f"幻灯片页数：{len(prs.slides)}")
