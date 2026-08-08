#!/usr/bin/env python3
"""生成 PPT：《东方修行路径对照表》讲解稿

青墨主题：墨青 + 青绿，避免常见紫渐变与奶油陶土风。
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

PRIMARY = RGBColor(0x1B, 0x3A, 0x4B)
ACCENT = RGBColor(0x2F, 0x6F, 0x5E)
LIGHT = RGBColor(0xE8, 0xF0, 0xED)
MIST = RGBColor(0xD5, 0xE5, 0xDE)
GOLD = RGBColor(0xA6, 0x7C, 0x3D)
CINNABAR = RGBColor(0x8B, 0x3A, 0x2F)
GREY = RGBColor(0x5A, 0x6A, 0x72)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x2A, 0x32)
BUDDHA = RGBColor(0x3A, 0x5A, 0x7A)
DAO = RGBColor(0x2F, 0x6F, 0x5E)
RU = RGBColor(0x7A, 0x5A, 0x3A)
BG_DEEP = RGBColor(0x12, 0x28, 0x34)


def set_font(run, name="微软雅黑", size=18, bold=False, color=DARK):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        rPr.append(rPr.makeelement(qn("a:ea"), {"typeface": name}))
    else:
        ea.set("typeface", name)


def add_rect(slide, left, top, width, height, fill=PRIMARY, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=18,
    bold=False,
    color=DARK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        text = [text]
    for i, t in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = t
        set_font(run, size=size, bold=bold, color=color)
    return tb


def add_bullets(slide, left, top, width, height, items, size=14, color=DARK, bullet="·"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.2
        run = p.add_run()
        run.text = f"{bullet}  {it}"
        set_font(run, size=size, color=color)
    return tb


def slide_chrome(slide, title, subtitle=None, page_no=None, total=None):
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.42), fill=PRIMARY)
    add_rect(slide, Inches(0), Inches(0.42), Inches(13.333), Inches(0.05), fill=ACCENT)
    add_text(
        slide,
        Inches(0.45),
        Inches(0.55),
        Inches(11),
        Inches(0.45),
        title,
        size=22,
        bold=True,
        color=PRIMARY,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    if subtitle:
        add_text(
            slide,
            Inches(0.45),
            Inches(0.95),
            Inches(12),
            Inches(0.3),
            subtitle,
            size=12,
            color=GREY,
            anchor=MSO_ANCHOR.MIDDLE,
        )
    if page_no and total:
        add_text(
            slide,
            Inches(11.6),
            Inches(7.1),
            Inches(1.4),
            Inches(0.25),
            f"{page_no} / {total}",
            size=10,
            color=GREY,
            align=PP_ALIGN.RIGHT,
        )


def make_card(slide, left, top, width, height, title, body_items, accent=ACCENT, title_size=14, body_size=12):
    add_rect(slide, left, top, width, height, fill=LIGHT)
    add_rect(slide, left, top, Inches(0.08), height, fill=accent)
    add_text(
        slide,
        left + Inches(0.2),
        top + Inches(0.12),
        width - Inches(0.3),
        Inches(0.35),
        title,
        size=title_size,
        bold=True,
        color=accent,
    )
    add_bullets(
        slide,
        left + Inches(0.2),
        top + Inches(0.5),
        width - Inches(0.3),
        height - Inches(0.6),
        body_items,
        size=body_size,
        color=DARK,
    )


def make_table(slide, left, top, width, height, headers, rows, header_color=PRIMARY, font_size=11):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    table = shape.table
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        cell.margin_left = Inches(0.06)
        cell.margin_right = Inches(0.06)
        cell.margin_top = Inches(0.04)
        cell.margin_bottom = Inches(0.04)
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = h
        set_font(r, size=font_size, bold=True, color=WHITE)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 == 0 else LIGHT
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            tf = cell.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(val)
            set_font(r, size=font_size - 1, color=DARK)
    return table


def path_block(slide, left, top, width, height, name, goal, method, barrier, pitfall, accent):
    add_rect(slide, left, top, width, height, fill=LIGHT)
    add_rect(slide, left, top, width, Inches(0.42), fill=accent)
    add_text(
        slide,
        left + Inches(0.15),
        top + Inches(0.05),
        width - Inches(0.3),
        Inches(0.35),
        name,
        size=16,
        bold=True,
        color=WHITE,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    lines = [
        f"目标：{goal}",
        f"方法：{method}",
        f"门槛：{barrier}",
        f"误区：{pitfall}",
    ]
    add_bullets(
        slide,
        left + Inches(0.15),
        top + Inches(0.55),
        width - Inches(0.25),
        height - Inches(0.65),
        lines,
        size=12,
        color=DARK,
    )


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    TOTAL = 14
    page = [0]

    def new_slide():
        page[0] += 1
        return prs.slides.add_slide(blank), page[0]

    # 1 封面
    s, _ = new_slide()
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), fill=BG_DEEP)
    # 氛围：斜向色带（抽象远山层次）
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(2.2), fill=PRIMARY)
    add_rect(s, Inches(0), Inches(5.6), Inches(13.333), Inches(1.9), fill=RGBColor(0x0E, 0x1F, 0x28))
    add_rect(s, Inches(0.8), Inches(2.55), Inches(0.1), Inches(2.0), fill=ACCENT)
    add_text(
        s,
        Inches(1.15),
        Inches(2.5),
        Inches(11),
        Inches(0.7),
        "东方修行路径对照表",
        size=40,
        bold=True,
        color=WHITE,
    )
    add_text(
        s,
        Inches(1.15),
        Inches(3.25),
        Inches(11),
        Inches(0.45),
        "佛家 · 道家 · 儒家",
        size=22,
        bold=True,
        color=RGBColor(0xA8, 0xD5, 0xC4),
    )
    add_text(
        s,
        Inches(1.15),
        Inches(3.85),
        Inches(11),
        Inches(0.4),
        "目标 · 方法 · 入门门槛 · 典型误区",
        size=16,
        color=MIST,
    )
    add_text(
        s,
        Inches(1.15),
        Inches(6.15),
        Inches(11),
        Inches(0.35),
        "知识对照 / 学习地图 ｜ V1.0 ｜ 非劝信、非替代师承",
        size=12,
        color=GREY,
    )

    # 2 目录
    s, p = new_slide()
    slide_chrome(s, "目录", "结构一览", p, TOTAL)
    cols = [
        (Inches(0.5), ["一、为何做对照表", "二、三教宏观对照", "三、佛家三条路径", "四、道家三条路径"]),
        (Inches(4.7), ["五、儒家两条路径", "六、四维横比总表", "七、入门路线图", "八、误区与边界"]),
        (Inches(8.9), ["附录取向", "先辨问题，再选方法", "一主一辅，忌混炖", "身心异常先就医"]),
    ]
    for left, items in cols[:2]:
        make_card(s, left, Inches(1.5), Inches(3.9), Inches(4.6), "内容", items, accent=ACCENT, body_size=15)
    make_card(s, cols[2][0], Inches(1.5), Inches(3.9), Inches(4.6), "使用原则", cols[2][1], accent=GOLD, body_size=15)

    # 3 为何对照
    s, p = new_slide()
    slide_chrome(s, "一、为何做对照表", "把「兴趣」翻译成「可选择的路径地图」", p, TOTAL)
    make_card(
        s,
        Inches(0.5),
        Inches(1.45),
        Inches(6.0),
        Inches(5.2),
        "常见困境",
        [
            "名词很多：禅、净土、丹道、慎独……不知从何起",
            "把不同问题的答案硬比高下（解脱 vs 成德）",
            "方法混炖：今天周天、明天话头、后天念佛",
            "忽视门槛：无师练高阶命功，或苛责式慎独",
            "被神通、速成、恐吓营销带偏",
        ],
        accent=CINNABAR,
    )
    make_card(
        s,
        Inches(6.8),
        Inches(1.45),
        Inches(6.0),
        Inches(5.2),
        "本表提供什么",
        [
            "按传统与路径拆开：目标 / 方法 / 门槛 / 误区",
            "三教总对照：先看问题意识差异",
            "情境化入门：忙人、爱静、爱读书、重身体……",
            "误区速查：从症状反查纠偏",
            "边界声明：对照理解，不替代传承与医疗",
        ],
        accent=ACCENT,
    )

    # 4 三教总对照
    s, p = new_slide()
    slide_chrome(s, "二、三教宏观对照", "问题不同，路径才不同", p, TOTAL)
    make_table(
        s,
        Inches(0.4),
        Inches(1.4),
        Inches(12.5),
        Inches(5.4),
        ["对照项", "佛家", "道家", "儒家"],
        [
            ["核心问题", "苦从何来？如何了生死？", "如何与道合真、全生？", "如何成德、安身立命？"],
            ["终极指向", "解脱 / 成佛", "合道 / 性命双修", "成圣贤 / 内圣外王"],
            ["方法气质", "戒定慧、信愿行、止观", "虚静、丹道、科仪、导引", "格致诚正、克己、礼"],
            ["社会位置", "出离与菩萨道并存", "隐逸与济世并存", "家国天下为主场"],
            ["入门印象", "净土易；禅中；唯识难", "导引易；符箓需传；丹难", "懂易行难，难在持续"],
        ],
        font_size=12,
    )

    # 5 佛家
    s, p = new_slide()
    slide_chrome(s, "三、佛家路径", "禅宗 · 净土 · 唯识", p, TOTAL)
    path_block(
        s,
        Inches(0.35),
        Inches(1.4),
        Inches(4.1),
        Inches(5.4),
        "禅宗",
        "见性成佛，日用保任",
        "坐禅、话头/公案、觉照",
        "中高，宜有师资",
        "口头禅、求玄、厌弃责任",
        BUDDHA,
    )
    path_block(
        s,
        Inches(4.6),
        Inches(1.4),
        Inches(4.1),
        Inches(5.4),
        "净土",
        "信愿持名，往生净土",
        "念佛、发愿、回向",
        "低，重在持续真切",
        "感应攀比、废事废戒",
        BUDDHA,
    )
    path_block(
        s,
        Inches(8.85),
        Inches(1.4),
        Inches(4.1),
        Inches(5.4),
        "唯识",
        "转识成智，教观双运",
        "经论 + 观心止观",
        "高，名相繁密",
        "知解障、炫博轻修",
        BUDDHA,
    )

    # 6 佛家要点
    s, p = new_slide()
    slide_chrome(s, "佛家共通地基", "三条路径不同，地基相似", p, TOTAL)
    make_card(s, Inches(0.5), Inches(1.45), Inches(4.0), Inches(5.2), "戒与伦理", ["不伤害、诚实、节制", "定慧需地基", "五戒精神先落地生活"], accent=BUDDHA)
    make_card(s, Inches(4.7), Inches(1.45), Inches(4.0), Inches(5.2), "正见与方向", ["缘起因果可转", "出离心 + 慈悲", "防神通导向与自欺"], accent=ACCENT)
    make_card(s, Inches(8.9), Inches(1.45), Inches(4.0), Inches(5.2), "善知识", ["可依止的人/团体", "多方了解再亲近", "远离恐吓与违法诱导"], accent=GOLD)

    # 7 道家
    s, p = new_slide()
    slide_chrome(s, "四、道家路径", "内丹 · 符箓 · 导引", p, TOTAL)
    path_block(
        s,
        Inches(0.35),
        Inches(1.4),
        Inches(4.1),
        Inches(5.4),
        "内丹",
        "性命双修，还虚合道",
        "筑基、调息、周天次第",
        "高，须明师",
        "盲练、气感崇拜、废德",
        DAO,
    )
    path_block(
        s,
        Inches(4.6),
        Inches(1.4),
        Inches(4.1),
        Inches(5.4),
        "符箓",
        "通真济世，科仪安顿",
        "授箓、斋醮、符咒、积德",
        "中，靠门派传承",
        "符力迷信、功利恐吓",
        DAO,
    )
    path_block(
        s,
        Inches(8.85),
        Inches(1.4),
        Inches(4.1),
        Inches(5.4),
        "导引",
        "形气神和，养生奠基",
        "导引、吐纳、站桩、动功",
        "低中，可渐进",
        "逞强过量、特效迷信",
        DAO,
    )

    # 8 道家提醒
    s, p = new_slide()
    slide_chrome(s, "道家实践提醒", "自然 · 节度 · 安全 · 传承", p, TOTAL)
    make_table(
        s,
        Inches(0.5),
        Inches(1.5),
        Inches(12.3),
        Inches(5.2),
        ["主题", "要义", "高危信号"],
        [
            ["自然与节度", "合节律，而非放纵或硬拼", "屏气硬顶、疼痛当进步"],
            ["德行为基", "无德求术易入歧途", "鼓励损人利己「用术」"],
            ["身体安全", "异常即停，医疗优先", "胸闷心悸仍强练、禁医"],
            ["传承辨别", "正统公开可核对", "高价速成、神秘恐吓、断亲"],
        ],
        header_color=DAO,
        font_size=13,
    )

    # 9 儒家
    s, p = new_slide()
    slide_chrome(s, "五、儒家路径", "修身 · 慎独——日用人伦即道场", p, TOTAL)
    path_block(
        s,
        Inches(0.5),
        Inches(1.4),
        Inches(6.0),
        Inches(5.4),
        "修身",
        "成德成人，推及家国",
        "格致诚正、克己复礼、事上磨",
        "低中：懂易，落在角色责任难",
        "道德表演、指责他人、空谈废事",
        RU,
    )
    path_block(
        s,
        Inches(6.8),
        Inches(1.4),
        Inches(6.0),
        Inches(5.4),
        "慎独",
        "独处亦正，意念不自欺",
        "省察克治、戒惧慎微、日省",
        "中：难在无人时的真实持续",
        "苛责成疾、形式反省、逃避责任",
        RU,
    )

    # 10 四维总表
    s, p = new_slide()
    slide_chrome(s, "六、四维横比（摘录）", "目标清晰度 × 门槛 × 风险的速览", p, TOTAL)
    make_table(
        s,
        Inches(0.35),
        Inches(1.35),
        Inches(12.6),
        Inches(5.6),
        ["路径", "近程目标", "综合门槛(1-5)", "师承依赖", "首要误区"],
        [
            ["禅宗", "心稳、少被动反应", "4", "强建议", "口头禅"],
            ["净土", "心有所归", "2", "可自学起步", "感应攀比"],
            ["唯识", "能辨心所", "5", "建议辅导", "知解障"],
            ["内丹", "精气神较充、欲有度", "5", "必须", "无师盲练"],
            ["符箓", "礼仪与服务安顿", "4", "必须", "符力迷信"],
            ["导引", "柔顺有力、息顺", "2", "教师更佳", "逞强过量"],
            ["修身", "一事能克己", "2", "不必秘传", "道德表演"],
            ["慎独", "独处少自欺", "3", "不必秘传", "苛责自我"],
        ],
        font_size=11,
    )

    # 11 入门路线
    s, p = new_slide()
    slide_chrome(s, "七、入门路线图", "按情境选择，一主一辅", p, TOTAL)
    make_table(
        s,
        Inches(0.35),
        Inches(1.35),
        Inches(12.6),
        Inches(5.6),
        ["情境", "建议主路径", "可搭配", "30天起步"],
        [
            ["忙、心不安", "净土持名", "儒家修身", "定课念佛 + 一件负责事"],
            ["喜安静坐", "禅宗坐禅", "导引松身", "静坐20分 + 松柔"],
            ["爱体系阅读", "唯识入门", "短时止观", "导读 + 观心日记"],
            ["身体紧、睡差", "导引/站桩", "作息节欲", "每日20分 + 固定睡眠"],
            ["有明师长线", "内丹依师", "导引+积德", "只做师嘱筑基"],
            ["重仪式服务", "正规符箓学习", "持诵积德", "先核传承与规戒"],
            ["关系中成长", "修身", "慎独", "读《大学》+事上克己"],
            ["易自欺", "慎独", "短静坐", "晚间三问日记"],
        ],
        font_size=11,
    )

    # 12 误区
    s, p = new_slide()
    slide_chrome(s, "八、误区与高危信号", "出现这些，停下来", p, TOTAL)
    make_card(
        s,
        Inches(0.45),
        Inches(1.4),
        Inches(6.1),
        Inches(5.3),
        "认知与心理误区",
        [
            "开悟自慢 / 口头禅 / 名相游戏",
            "感应攀比、恐吓式信仰",
            "道德表演、以儒学压人",
            "苛责式慎独导致自我厌恶",
            "神通导向、境界炫耀",
        ],
        accent=CINNABAR,
        body_size=14,
    )
    make_card(
        s,
        Inches(6.8),
        Inches(1.4),
        Inches(6.1),
        Inches(5.3),
        "应立即远离 / 就医的信号",
        [
            "胸闷心悸仍强练；禁止就医",
            "高价速成密法 + 神秘恐吓",
            "鼓励违法、毁家庭、全面断亲",
            "持续抑郁焦虑却禁止求助",
            "以符咒/修行替代急症医疗",
        ],
        accent=GOLD,
        body_size=14,
    )

    # 13 组合原则
    s, p = new_slide()
    slide_chrome(s, "实践组合原则", "现代生活中的稳健用法", p, TOTAL)
    principles = [
        ("一主一辅", "同时只深耕一条主路径，另一条滋养，勿三线深钻"),
        ("身→心→理", "身躁先导引作息；心乱立定课；理路后置"),
        ("入世不废", "有家庭职场者，优先能嵌入日程的路径"),
        ("高门槛单列", "内丹/符箓/深度禅堂，单独评估师承与身心"),
        ("季度复盘", "更平静？更负责？更少自欺？无效则调整"),
        ("对照非混一", "可互相启发，不可把术语体系强行合一"),
    ]
    for i, (t, b) in enumerate(principles):
        col = i % 3
        row = i // 3
        left = Inches(0.45 + col * 4.25)
        top = Inches(1.45 + row * 2.7)
        make_card(s, left, top, Inches(4.05), Inches(2.45), t, [b], accent=ACCENT if i % 2 == 0 else PRIMARY, body_size=13)

    # 14 收束
    s, p = new_slide()
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), fill=BG_DEEP)
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.5), fill=ACCENT)
    add_text(
        s,
        Inches(1.0),
        Inches(2.2),
        Inches(11.3),
        Inches(0.7),
        "先问自己在解什么题",
        size=32,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        Inches(1.5),
        Inches(3.2),
        Inches(10.3),
        Inches(1.2),
        [
            "苦与生死 → 佛家　　合道与形神 → 道家　　成人与人伦 → 儒家",
            "选一条主路，守住门槛与误区，让日用成为道场。",
        ],
        size=16,
        color=MIST,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        Inches(1.0),
        Inches(5.5),
        Inches(11.3),
        Inches(0.4),
        "配套交付：东方修行路径对照表.xlsx（九个工作表详表）",
        size=14,
        color=GOLD,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        Inches(1.0),
        Inches(6.2),
        Inches(11.3),
        Inches(0.35),
        f"{p} / {TOTAL}",
        size=11,
        color=GREY,
        align=PP_ALIGN.CENTER,
    )

    out_dir = Path(__file__).resolve().parents[1] / "deliverables"
    out_dir.mkdir(parents=True, exist_ok=True)
    out = out_dir / "东方修行路径对照表.pptx"
    prs.save(out)
    print(f"已生成: {out}")


if __name__ == "__main__":
    main()
