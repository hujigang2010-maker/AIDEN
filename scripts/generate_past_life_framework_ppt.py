# -*- coding: utf-8 -*-
"""生成《前世叙事的功能分析框架》PPT。

四维功能：心理补偿 · 关系解释 · 身份建构 · 道德约束
并对照不同文化传统中前世故事各自承担的功能重心。
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ----------------------------------------------------------------------------
# 视觉方向：夜间书斋 —— 墨蓝底 + 青瓷强调 + 暖琥珀点缀
# （避开紫白渐变 / 奶油赤陶 / 报纸排版默认风）
# ----------------------------------------------------------------------------
INK = RGBColor(0x10, 0x14, 0x1C)
INK_2 = RGBColor(0x18, 0x1E, 0x2A)
PANEL = RGBColor(0x22, 0x2A, 0x3A)
PANEL_LT = RGBColor(0x2C, 0x36, 0x48)
CELADON = RGBColor(0x6B, 0xA8, 0x96)  # 青瓷
CELADON_DK = RGBColor(0x3D, 0x6E, 0x62)
AMBER = RGBColor(0xC9, 0x96, 0x5A)  # 琥珀
AMBER_LT = RGBColor(0xE2, 0xC0, 0x8A)
WHITE = RGBColor(0xF4, 0xF1, 0xEA)
OFF = RGBColor(0xD8, 0xD4, 0xCB)
MUTED = RGBColor(0x8E, 0x96, 0xA4)
CORAL = RGBColor(0xC0, 0x6A, 0x5A)  # 朱砂偏冷，用于警示/张力

FONT = "微软雅黑"
SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]
PAGE = 0


def _font(run, size, color, bold=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = FONT
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT)


def slide(bg=INK):
    global PAGE
    PAGE += 1
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid()
    r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def rect(s, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE, line=None):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line:
        sp.line.color.rgb = line
        sp.line.width = Pt(1.25)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def rrect(s, x, y, w, h, color, line=None, adj=0.08):
    sp = rect(s, x, y, w, h, color, MSO_SHAPE.ROUNDED_RECTANGLE, line)
    try:
        sp.adjustments[0] = adj
    except Exception:
        pass
    return sp


def txt(s, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, size, color, bold[, space_after])"""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    for i, spec in enumerate(lines):
        text, size, color, bold = spec[0], spec[1], spec[2], spec[3]
        space_after = spec[4] if len(spec) > 4 else 4
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = text
        _font(run, size, color, bold)
    return tb


def footer(s, page_label=None):
    label = page_label or f"{PAGE:02d}"
    txt(
        s,
        Inches(0.5),
        Inches(7.05),
        Inches(10),
        Inches(0.3),
        [("前世叙事功能分析框架  ·  AIDEN", 10, MUTED, False)],
    )
    txt(
        s,
        Inches(11.6),
        Inches(7.05),
        Inches(1.2),
        Inches(0.3),
        [(label, 10, MUTED, False)],
        align=PP_ALIGN.RIGHT,
    )


def header(s, title, subtitle=None):
    rect(s, Inches(0), Inches(0), Inches(0.14), SH, CELADON)
    txt(s, Inches(0.5), Inches(0.28), Inches(12), Inches(0.5), [(title, 24, WHITE, True)])
    if subtitle:
        txt(
            s,
            Inches(0.52),
            Inches(0.78),
            Inches(12),
            Inches(0.35),
            [(subtitle, 12, MUTED, False)],
        )



# ============================================================================ ============================================================================icons
# 封面
# ============================================================================iconsicons
s = slide()
# atmospheric bands
rect(s, 0, 0, SW, Inches(0.08), CELADON)
rect(s, 0, Inches(7.42), SW, Inches(0.08), AMBER)
txt(
    s,
    Inches(0.9),
    Inches(1.6),
    Inches(11.5),
    Inches(0.4),
    [("AI 可用知识框架  ·  叙事功能学", 14, CELADON, False)],
)
txt(
    s,
    Inches(0.9),
    Inches(2.15),
    Inches(11.5),
    Inches(1.2),
    [("前世叙事的功能分析框架", 40, WHITE, True, 8)],
)
txt(
    s,
    Inches(0.9),
    Inches(3.4),
    Inches(11),
    Inches(0.8),
    [
        (
            "心理补偿 · 关系解释 · 身份建构 · 道德约束",
            18,
            AMBER_LT,
            False,
            10,
        ),
        (
            "不同文化里的前世故事，分别在承担什么功能？",
            16,
            OFF,
            False,
        ),
    ],
)
txt(
    s,
    Inches(0.9),
    Inches(5.6),
    Inches(10),
    Inches(0.6),
    [
        ("供内容分析、标注体系与 AI 提示词沉淀使用", 13, MUTED, False),
        ("v1.0", 12, MUTED, False),
    ],
)

# ============================================================================iconsicons
# 目录 / 为何需要
# ============================================================================icons ============================================================================
s = slide()
header(s, "为何需要「功能分析」", "把前世故事当作社会—心理装置，而不是只问「信不信」")
footer(s)

points = [
    ("01", "现象普遍", "从宗教经典、民间传说到短视频「前世今生」，叙事持续再生产。"),
    ("02", "功能稳定", "具体情节多变，但服务的心理与社会功能高度收敛。"),
    ("03", "文化分流", "同一功能在不同传统中权重不同，载体与合法性来源也不同。"),
    ("04", "AI 可操作", "可转成标注维度、检索标签、生成约束与跨文化对照提示。"),
]
for i, (n, t, d) in enumerate(points):
    y = Inches(1.35) + Inches(i * 1.25)
    rrect(s, Inches(0.55), y, Inches(12.2), Inches(1.1), PANEL, adj=0.06)
    rrect(s, Inches(0.8), y + Inches(0.28), Inches(0.7), Inches(0.55), CELADON_DK, adj=0.15)
    txt(
        s,
        Inches(0.8),
        y + Inches(0.35),
        Inches(0.7),
        Inches(0.45),
        [(n, 14, WHITE, True)],
        align=PP_ALIGN.CENTER,
    )
    txt(s, Inches(1.8), y + Inches(0.22), Inches(10.5), Inches(0.4), [(t, 18, WHITE, True)])
    txt(s, Inches(1.8), y + Inches(0.58), Inches(10.5), Inches(0.4), [(d, 13, OFF, False)])

# ============================================================================iconsicons
# 总览四维
# ============================================================================iconsicons
s = slide()
header(s, "四维功能总览", "Past-Life Narrative Functional Framework（PLN-F）")
footer(s)

dims = [
    ("F1", "心理补偿", "Psychological\nCompensation", "用「另一世的我」消化创伤、遗憾、不公与未完成之事", CELADON),
    ("F2", "关系解释", "Relational\nExplanation", "为亲密、疏离、冲突与缘分提供超验因果剧本", AMBER),
    ("F3", "身份建构", "Identity\nConstruction", "借前世地位/天赋/使命重写当下的自我叙事", RGBColor(0x7A, 0x9E, 0xC2)),
    ("F4", "道德约束", "Moral\nConstraint", "以跨世因果规范行为，把伦理后果拉长到来世", CORAL),
]
for i, (code, zh, en, desc, color) in enumerate(dims):
    x = Inches(0.45) + Inches(i * 3.2)
    rrect(s, x, Inches(1.4), Inches(3.0), Inches(5.1), PANEL, adj=0.05)
    rect(s, x, Inches(1.4), Inches(3.0), Inches(0.12), color)
    txt(s, x + Inches(0.2), Inches(1.75), Inches(2.6), Inches(0.35), [(code, 12, color, True)])
    txt(s, x + Inches(0.2), Inches(2.2), Inches(2.6), Inches(0.5), [(zh, 22, WHITE, True)])
    for j, line in enumerate(en.split("\n")):
        txt(
            s,
            x + Inches(0.2),
            Inches(2.8) + Inches(j * 0.28),
            Inches(2.6),
            Inches(0.3),
            [(line, 11, MUTED, False)],
        )
    txt(s, x + Inches(0.2), Inches(3.7), Inches(2.6), Inches(2.2), [(desc, 13, OFF, False, 8)])

# ============================================================================iconsicons
# F1 心理补偿
# ============================================================================icons========
s = slide()
header(s, "F1  心理补偿", "用叙事把不可承受的偶然，改写成可理解的「早有安排」")
footer(s)

rrect(s, Inches(0.5), Inches(1.35), Inches(6.0), Inches(5.2), PANEL, adj=0.05)
txt(s, Inches(0.75), Inches(1.55), Inches(5.5), Inches(0.4), [("典型触发情境", 16, CELADON, True)])
for i, line in enumerate(
    [
        "早逝、残疾、重大丧失 → 「前世欠债 / 功课未完」",
        "阶层固化与人生挫败 → 「前世曾显赫 / 今世修行」",
        "创伤记忆碎片 → 催眠回溯中的「前世场景」",
        "疾病与痛苦 → 「业障消尽」式意义化",
        "对平凡自我的不满 → 「沉睡的特殊灵魂」",
    ]
):
    txt(
        s,
        Inches(0.75),
        Inches(2.15) + Inches(i * 0.7),
        Inches(5.5),
        Inches(0.6),
        [(line, 13, OFF, False)],
    )

rrect(s, Inches(6.8), Inches(1.35), Inches(5.95), Inches(5.2), PANEL, adj=0.05)
txt(s, Inches(7.05), Inches(1.55), Inches(5.5), Inches(0.4), [("功能机制与识别信号", 16, AMBER, True)])
mech = [
    ("意义修复", "把随机苦难纳入目的论秩序"),
    ("责任外移", "部分卸下当下自我的全责"),
    ("希望延期", "「来世可偿」缓冲绝望"),
    ("识别信号", "疗愈话术、遗憾闭环、不公合理化、回溯治疗叙事"),
]
for i, (t, d) in enumerate(mech):
    y = Inches(2.2) + Inches(i * 0.95)
    txt(s, Inches(7.05), y, Inches(5.4), Inches(0.35), [(t, 14, WHITE, True)])
    txt(s, Inches(7.05), y + Inches(0.35), Inches(5.4), Inches(0.4), [(d, 12, MUTED, False)])

# ============================================================================icons========
# F2 关系解释
# ============================================================================icons========
s = slide()
header(s, "F2  关系解释", "前世是人际关系的「超验说明书」")
footer(s)

cards = [
    ("吸引与命定", "一见钟情、灵魂伴侣、双生火焰 → 前世未了情", CELADON),
    ("冲突与纠缠", "反复伤害的关系 → 债主—债务人、仇敌转世", CORAL),
    ("亲缘重组", "子女是「讨债/还债」来的；养父母是前世亲人", AMBER),
    ("社群边界", "同门、同修、同部落灵魂 → 圈内信任的神话基础", RGBColor(0x7A, 0x9E, 0xC2)),
]
for i, (t, d, c) in enumerate(cards):
    col, row = i % 2, i // 2
    x = Inches(0.5) + Inches(col * 6.4)
    y = Inches(1.4) + Inches(row * 2.55)
    rrect(s, x, y, Inches(6.1), Inches(2.3), PANEL, adj=0.05)
    rect(s, x, y, Inches(0.12), Inches(2.3), c)
    txt(s, x + Inches(0.4), y + Inches(0.4), Inches(5.4), Inches(0.45), [(t, 18, WHITE, True)])
    txt(s, x + Inches(0.4), y + Inches(1.05), Inches(5.4), Inches(0.9), [(d, 14, OFF, False)])

# ============================================================================icons========
# F3 身份建构
# ============================================================================icons========
s = slide()
header(s, "F3  身份建构", "前世是自我的「升级补丁」与合法性来源")
footer(s)

left_items = [
    ("贵族/圣者前世", "抬升当下地位，对抗边缘感"),
    ("特殊使命叙事", "「我此生是来完成某任务」"),
    ("天赋溯源", "艺术、通灵、领导力「与生俱来」"),
    ("转世制度身份", "活佛、祖灵回归等制度化认证"),
]
rrect(s, Inches(0.5), Inches(1.35), Inches(6.1), Inches(5.2), PANEL, adj=0.05)
txt(s, Inches(0.75), Inches(1.55), Inches(5.6), Inches(0.4), [("建构路径", 16, CELADON, True)])
for i, (t, d) in enumerate(left_items):
    y = Inches(2.2) + Inches(i * 0.95)
    txt(s, Inches(0.75), y, Inches(5.5), Inches(0.35), [(t, 15, WHITE, True)])
    txt(s, Inches(0.75), y + Inches(0.38), Inches(5.5), Inches(0.4), [(d, 13, MUTED, False)])

rrect(s, Inches(6.9), Inches(1.35), Inches(5.85), Inches(5.2), PANEL, adj=0.05)
txt(s, Inches(7.15), Inches(1.55), Inches(5.4), Inches(0.4), [("张力与风险", 16, CORAL, True)])
risks = [
    "精英幻想：人人皆法老/王妃的通胀叙事",
    "权威挪用：以「我前世是…」压过对话平等",
    "历史挪用：对真实族群/人物的浪漫化占有",
    "制度版 vs 民间版：认证成本决定可信度",
    "AI 生成风险：批量制造虚假「高贵前世」",
]
for i, line in enumerate(risks):
    txt(
        s,
        Inches(7.15),
        Inches(2.2) + Inches(i * 0.75),
        Inches(5.4),
        Inches(0.65),
        [(f"·  {line}", 13, OFF, False)],
    )

# ============================================================================icons========
# F4 道德约束
# ============================================================================icons========
s = slide()
header(s, "F4  道德约束", "把伦理账本延伸到「看不见的来世」")
footer(s)

rows = [
    ("因果报应", "善恶有报被时间轴拉长，弥补当世不公的可见性缺口"),
    ("戒律强化", "杀生、妄语、破戒 → 来世堕恶趣；布施持戒 → 人天福报"),
    ("社会秩序", "劝孝、劝善、止争；用恐惧与希望双轨维稳日常伦理"),
    ("自我规训", "个体以内疚/功德心管理欲望，降低外在监控成本"),
    ("识别信号", "「造业」「消业」「功德」「恶报」等跨世奖惩词汇簇"),
]
for i, (t, d) in enumerate(rows):
    y = Inches(1.35) + Inches(i * 1.0)
    rrect(s, Inches(0.5), y, Inches(12.3), Inches(0.88), PANEL, adj=0.06)
    rrect(s, Inches(0.7), y + Inches(0.2), Inches(2.2), Inches(0.48), CELADON_DK, adj=0.15)
    txt(
        s,
        Inches(0.7),
        y + Inches(0.28),
        Inches(2.2),
        Inches(0.4),
        [(t, 13, WHITE, True)],
        align=PP_ALIGN.CENTER,
    )
    txt(s, Inches(3.2), y + Inches(0.28), Inches(9.2), Inches(0.5), [(d, 14, OFF, False)])

# ============================================================================icons========
# 跨文化矩阵
# ============================================================================icons========
s = slide()
header(s, "跨文化功能矩阵", "权重示意：● 强  ◐ 中  ○ 弱（分析用，非绝对）")
footer(s)

# table header
cultures = ["印度教/佛教", "汉传民间", "藏传转世", "日本流行", "西方新时代", "非洲祖灵"]
funcs = ["心理补偿", "关系解释", "身份建构", "道德约束"]
# weights matrix
matrix = [
    ["◐", "◐", "●", "●"],  # Indic
    ["◐", "●", "○", "●"],  # Chinese folk
    ["○", "◐", "●", "◐"],  # Tibetan
    ["◐", "●", "◐", "○"],  # Japanese pop
    ["●", "●", "●", "○"],  # New Age
    ["◐", "◐", "●", "◐"],  # African ancestor
]

# header row
rrect(s, Inches(0.45), Inches(1.35), Inches(12.4), Inches(0.55), PANEL_LT, adj=0.04)
txt(s, Inches(0.55), Inches(1.45), Inches(2.4), Inches(0.4), [("文化传统", 12, MUTED, True)])
for j, f in enumerate(funcs):
    txt(
        s,
        Inches(3.1) + Inches(j * 2.4),
        Inches(1.45),
        Inches(2.2),
        Inches(0.4),
        [(f, 12, CELADON, True)],
        align=PP_ALIGN.CENTER,
    )

for i, cult in enumerate(cultures):
    y = Inches(2.0) + Inches(i * 0.75)
    bg = PANEL if i % 2 == 0 else INK_2
    rrect(s, Inches(0.45), y, Inches(12.4), Inches(0.68), bg, adj=0.04)
    txt(s, Inches(0.6), y + Inches(0.15), Inches(2.4), Inches(0.4), [(cult, 13, WHITE, True)])
    for j, w in enumerate(matrix[i]):
        color = AMBER if w == "●" else (OFF if w == "◐" else MUTED)
        txt(
            s,
            Inches(3.1) + Inches(j * 2.4),
            y + Inches(0.15),
            Inches(2.2),
            Inches(0.4),
            [(w, 16, color, True)],
            align=PP_ALIGN.CENTER,
        )

# ============================================================================icons========
# 文化案例 1
# ====================================================================================
s = slide()
header(s, "文化切片（一）", "印度—佛教传统 · 汉传民间因果")
footer(s)

rrect(s, Inches(0.5), Inches(1.35), Inches(6.1), Inches(5.2), PANEL, adj=0.05)
txt(s, Inches(0.75), Inches(1.55), Inches(5.5), Inches(0.4), [("印度教 / 佛教轮回", 16, CELADON, True)])
for line in [
    "核心装置：业（karma）—轮回（saṃsāra）—解脱",
    "F4 极强：伦理后果跨世兑现，支撑戒律体系",
    "F3 强：种姓、阿罗汉、菩萨道等身份阶序",
    "F1/F2 存在但常被教理化，非私人疗愈口语",
    "AI 提示：区分教理「业果」与民间「报应故事」",
]:
    pass
bullets = [
    "核心装置：业—轮回—解脱",
    "F4 极强：伦理后果跨世兑现",
    "F3 强：修行阶位与种姓叙事交织",
    "F1/F2 存在，但常被教理化",
    "分析要点：教理文本 vs 民间报应话本",
]
for i, line in enumerate(bullets):
    txt(
        s,
        Inches(0.75),
        Inches(2.2) + Inches(i * 0.7),
        Inches(5.5),
        Inches(0.6),
        [(f"·  {line}", 13, OFF, False)],
    )

rrect(s, Inches(6.9), Inches(1.35), Inches(5.85), Inches(5.2), PANEL, adj=0.05)
txt(s, Inches(7.15), Inches(1.55), Inches(5.4), Inches(0.4), [("汉传民间 · 缘与报", 16, AMBER, True)])
bullets2 = [
    "核心装置：缘分、投胎、因果报应、阴德",
    "F2 极强：「前世冤家」「今世夫妻」",
    "F4 强：善书、宝卷、劝善戏曲",
    "F3 弱于制度化转世，多个人传说",
    "分析要点：家庭伦理剧是主要载体",
]
for i, line in enumerate(bullets2):
    txt(
        s,
        Inches(7.15),
        Inches(2.2) + Inches(i * 0.7),
        Inches(5.4),
        Inches(0.6),
        [(f"·  {line}", 13, OFF, False)],
    )

# ====================================================================================
# 文化案例 2
# ====================================================================================
s = slide()
header(s, "文化切片（二）", "藏传转世制度 · 日本流行文化 · 西方新时代 · 非洲祖灵")
footer(s)

cases = [
    ("藏传转世（活佛）", "F3 制度化巅峰：认证、寻访、坐床；身份即政治—宗教资本", CELADON),
    ("日本流行前世", "动漫/小说中 F2 主导：命定重逢、轮回恋人；道德弱化", AMBER),
    ("西方新时代", "催眠回溯、灵魂伴侣课：F1+F2+F3 消费化；F4 淡化", RGBColor(0x7A, 0x9E, 0xC2)),
    ("非洲祖灵回归", "祖先在氏族中再生：F3 社群连续 + 伦理义务，非个人幻想", CORAL),
]
for i, (t, d, c) in enumerate(cases):
    col, row = i % 2, i // 2
    x = Inches(0.5) + Inches(col * 6.4)
    y = Inches(1.4) + Inches(row * 2.55)
    rrect(s, x, y, Inches(6.1), Inches(2.3), PANEL, adj=0.05)
    rect(s, x, y, Inches(0.12), Inches(2.3), c)
    txt(s, x + Inches(0.4), y + Inches(0.45), Inches(5.4), Inches(0.45), [(t, 17, WHITE, True)])
    txt(s, x + Inches(0.4), y + Inches(1.1), Inches(5.4), Inches(0.9), [(d, 13, OFF, False)])

# ====================================================================================
# 功能互动
# ====================================================================================
s = slide()
header(s, "四维如何联动", "真实文本很少单功能；分析要标「主功能 + 辅功能」")
footer(s)

flows = [
    ("补偿 → 关系", "疗愈叙事常落脚于「前世爱人出现」，F1 借 F2 完成闭环"),
    ("身份 → 道德", "「我曾是修行者」既抬升自我，也施加戒律式自我期待"),
    ("道德 → 补偿", "「受苦是消业」把惩罚语言改写成安慰语言"),
    ("关系 → 身份", "「我们是灵魂家族」同时划定圈子与抬高圈内人"),
]
for i, (t, d) in enumerate(flows):
    y = Inches(1.4) + Inches(i * 1.25)
    rrect(s, Inches(0.55), y, Inches(12.2), Inches(1.1), PANEL, adj=0.06)
    rrect(s, Inches(0.8), y + Inches(0.28), Inches(2.8), Inches(0.55), CELADON_DK, adj=0.12)
    txt(
        s,
        Inches(0.8),
        y + Inches(0.38),
        Inches(2.8),
        Inches(0.4),
        [(t, 14, WHITE, True)],
        align=PP_ALIGN.CENTER,
    )
    txt(s, Inches(3.9), y + Inches(0.35), Inches(8.5), Inches(0.5), [(d, 14, OFF, False)])

# ====================================================================================
# AI 标注 schema
# ====================================================================================
s = slide()
header(s, "AI 标注与分析协议", "把框架落成可执行的字段，而非只停留在概念")
footer(s)

fields = [
    ("primary_function", "主功能：F1–F4 单选"),
    ("secondary_functions", "辅功能：多选，可空"),
    ("culture_frame", "文化框架标签（教理/民间/流行/治疗…）"),
    ("narrative_device", "装置词：业、缘、债、使命、回溯…"),
    ("beneficiary", "功能服务对象：自我 / 关系双方 / 社群 / 制度"),
    ("legitimacy_claim", "合法性来源：经典、大师、催眠、梦、族谱…"),
    ("risk_flags", "风险：历史挪用、操控话术、贩卖焦虑等"),
]
for i, (k, v) in enumerate(fields):
    y = Inches(1.3) + Inches(i * 0.72)
    rrect(s, Inches(0.5), y, Inches(12.3), Inches(0.62), PANEL, adj=0.05)
    txt(s, Inches(0.75), y + Inches(0.15), Inches(3.8), Inches(0.4), [(k, 13, CELADON, True)])
    txt(s, Inches(4.7), y + Inches(0.15), Inches(7.8), Inches(0.4), [(v, 13, OFF, False)])

# ====================================================================================
# 提示词模板
# ====================================================================================
s = slide()
header(s, "分析提示词模板（可直接给模型）", "输入一段前世叙事文本，输出结构化功能画像")
footer(s)

prompt_lines = [
    "请用 PLN-F 四维框架分析下列前世叙事：",
    "1) 判定 primary_function 与 secondary_functions，并引用原文依据；",
    "2) 指出 narrative_device 与 culture_frame；",
    "3) 说明该叙事对谁产生功能收益（beneficiary）；",
    "4) 若跨文化改写，功能权重可能如何漂移；",
    "5) 标出潜在 risk_flags（操控、历史挪用、贩卖焦虑等）。",
    "输出 JSON，字段与框架 schema 对齐；结论须可回溯到文本。",
]
rrect(s, Inches(0.55), Inches(1.35), Inches(12.2), Inches(5.2), PANEL, adj=0.05)
for i, line in enumerate(prompt_lines):
    color = AMBER_LT if i == 0 else OFF
    bold = i == 0
    txt(
        s,
        Inches(0.9),
        Inches(1.65) + Inches(i * 0.6),
        Inches(11.5),
        Inches(0.55),
        [(line, 14, color, bold)],
    )

# ====================================================================================
# 应用场景
# ====================================================================================
s = slide()
header(s, "应用场景", "内容研究 · 产品设计 · 风险治理")
footer(s)

apps = [
    ("内容研究", "短视频/小说/宗教传播中的前世母题量化与跨文化对比"),
    ("咨询与疗愈伦理", "区分安慰功能与操控风险，设计知情同意话术"),
    ("叙事生成约束", "AI 写「前世故事」时控制功能配比，避免精英前世通胀"),
    ("标注训练数据", "为分类模型提供稳定标签体系与冲突裁定规则"),
    ("文化敏感性审查", "识别对真实宗教制度/族群历史的轻率挪用"),
    ("产品知识库", "作为 RAG 片段：解释「这段内容在干什么」而非「真不真」"),
]
for i, (t, d) in enumerate(apps):
    col, row = i % 3, i // 3
    x = Inches(0.45) + Inches(col * 4.25)
    y = Inches(1.4) + Inches(row * 2.55)
    rrect(s, x, y, Inches(4.05), Inches(2.3), PANEL, adj=0.05)
    rect(s, x, y, Inches(4.05), Inches(0.1), CELADON if row == 0 else AMBER)
    txt(s, x + Inches(0.25), y + Inches(0.4), Inches(3.55), Inches(0.45), [(t, 16, WHITE, True)])
    txt(s, x + Inches(0.25), y + Inches(1.0), Inches(3.55), Inches(1.0), [(d, 13, OFF, False)])

# ====================================================================================
# 使用边界
# ====================================================================================
s = slide()
header(s, "使用边界与学术姿态", "功能分析 ≠ 真伪裁判；尊重传统内部的真理诉求")
footer(s)

bounds = [
    ("不做的事", "不宣称前世「存在/不存在」；不把信仰降格为「纯幻觉」"),
    ("做的事", "描述叙事在心理与社会层面「完成了什么工作」"),
    ("文化谦逊", "矩阵是启发式工具；深入个案需结合民族志与教义学"),
    ("伦理红线", "不以框架协助操控性营销或针对脆弱人群的恐惧贩卖"),
    ("迭代方式", "用新语料回标，校正权重与装置词典，保持版本号"),
]
for i, (t, d) in enumerate(bounds):
    y = Inches(1.35) + Inches(i * 1.0)
    rrect(s, Inches(0.5), y, Inches(12.3), Inches(0.88), PANEL, adj=0.06)
    txt(s, Inches(0.8), y + Inches(0.25), Inches(2.4), Inches(0.45), [(t, 14, AMBER, True)])
    txt(s, Inches(3.4), y + Inches(0.25), Inches(9.0), Inches(0.5), [(d, 14, OFF, False)])

# ====================================================================================
# 收束
# ====================================================================================
s = slide()
rect(s, 0, 0, SW, Inches(0.08), CELADON)
rect(s, 0, Inches(7.42), SW, Inches(0.08), AMBER)
txt(
    s,
    Inches(0.9),
    Inches(2.0),
    Inches(11.5),
    Inches(0.5),
    [("一句话收束", 14, CELADON, False)],
)
txt(
    s,
    Inches(0.9),
    Inches(2.6),
    Inches(11.5),
    Inches(1.5),
    [
        (
            "前世叙事的核心价值，往往不在「证明另一世」，",
            26,
            WHITE,
            True,
            12,
        ),
        (
            "而在「安放这一世」——补偿、解释关系、建构身份、约束道德。",
            26,
            WHITE,
            True,
        ),
    ],
)
txt(
    s,
    Inches(0.9),
    Inches(5.0),
    Inches(11.5),
    Inches(0.8),
    [
        ("配套交付：framework/PLN-F.md  ·  framework/pln_f.schema.json", 13, MUTED, False),
        ("脚本：scripts/generate_past_life_framework_ppt.py", 13, MUTED, False),
    ],
)

# save
out_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "output")
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, "前世叙事的功能分析框架.pptx")
prs.save(out_path)
print(f"已生成：{out_path}  共 {PAGE} 页")
