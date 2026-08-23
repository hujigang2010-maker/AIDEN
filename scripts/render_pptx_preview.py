#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""把 PPTX 中的形状与文字栅格化成 PNG，便于走查。"""

from __future__ import annotations

import os
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "output" / "七款中文AI-PPT-Skill横评.pptx"
OUT = ROOT / "output" / "preview"
FONT_PATH = "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc"
SCALE = 1920 / 13.333  # px per inch


def emu_px(emu) -> int:
    inches = int(emu) / 914400
    return int(round(inches * SCALE))


def rgb_of(color) -> tuple[int, int, int] | None:
    if color is None:
        return None
    try:
        c = color.rgb
        return (c[0], c[1], c[2]) if not hasattr(c, "r") else (c[0], c[1], c[2])
    except Exception:
        return None


def fill_rgb(shape) -> tuple[int, int, int] | None:
    try:
        fill = shape.fill
        ftype = str(fill.type)
        if "BACKGROUND" in ftype or fill.type is None:
            return None
        return rgb_of(fill.fore_color)
    except Exception:
        return None


def bg_rgb(slide) -> tuple[int, int, int]:
    try:
        return rgb_of(slide.background.fill.fore_color) or (244, 240, 230)
    except Exception:
        return (244, 240, 230)


def font(size_pt: float, bold=False):
    px = max(10, int(size_pt * SCALE / 96 * 96 / 72 * 0.92))
    # pt -> px at 144 dpi-ish: 1pt = SCALE/72 px
    px = max(10, int(size_pt * SCALE / 72))
    try:
        return ImageFont.truetype(FONT_PATH, px)
    except Exception:
        return ImageFont.load_default()


def wrap(draw, text, fnt, max_w):
    if not text:
        return [""]
    lines = []
    for para in text.split("\n"):
        cur = ""
        for ch in para:
            trial = cur + ch
            if draw.textlength(trial, font=fnt) <= max_w or not cur:
                cur = trial
            else:
                lines.append(cur)
                cur = ch
        lines.append(cur)
    return lines or [""]


def draw_text_frame(img, shape):
    draw = ImageDraw.Draw(img)
    x, y = emu_px(shape.left), emu_px(shape.top)
    w, h = max(1, emu_px(shape.width)), max(1, emu_px(shape.height))
    tf = shape.text_frame
    # vertical anchor approximation
    total_h = 0
    blocks = []
    for p in tf.paragraphs:
        text = "".join(r.text or "" for r in p.runs) or p.text
        if not p.runs and not text:
            continue
        size = 14
        fill = (17, 17, 17)
        bold = False
        if p.runs:
            r0 = p.runs[0]
            if r0.font.size:
                size = r0.font.size.pt
            try:
                rgb = r0.font.color.rgb
                fill = (int(str(rgb)[0:2], 16), int(str(rgb)[2:4], 16), int(str(rgb)[4:6], 16)) if False else (
                    rgb[0], rgb[1], rgb[2]
                )
            except Exception:
                pass
            bold = bool(r0.font.bold)
        align = str(p.alignment) if p.alignment else "LEFT"
        fnt = font(size, bold)
        lines = wrap(draw, text, fnt, w - 4)
        line_h = int(size * SCALE / 72 * 1.18)
        blocks.append((lines, fnt, fill, align, line_h))
        total_h += line_h * len(lines)
    cy = y
    try:
        va = str(tf.vertical_anchor)
        if "MIDDLE" in va:
            cy = y + max(0, (h - total_h) // 2)
        elif "BOTTOM" in va:
            cy = y + max(0, h - total_h)
    except Exception:
        pass
    for lines, fnt, fill, align, line_h in blocks:
        for line in lines:
            tw = draw.textlength(line, font=fnt)
            if "RIGHT" in align:
                tx = x + w - tw
            elif "CENTER" in align:
                tx = x + (w - tw) / 2
            else:
                tx = x
            draw.text((tx, cy), line, font=fnt, fill=fill)
            cy += line_h


def render_slide(slide, path: Path):
    W, H = int(13.333 * SCALE), int(7.5 * SCALE)
    img = Image.new("RGB", (W, H), bg_rgb(slide))
    draw = ImageDraw.Draw(img)
    for shape in slide.shapes:
        x, y = emu_px(shape.left), emu_px(shape.top)
        w, h = max(1, emu_px(shape.width)), max(1, emu_px(shape.height))
        color = fill_rgb(shape)
        if color and w > 0 and h > 0:
            draw.rectangle([x, y, x + w, y + h], fill=color)
        if shape.has_text_frame:
            draw_text_frame(img, shape)
    img.save(path, "PNG", optimize=True)


def main():
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(PPTX))
    paths = []
    for i, slide in enumerate(prs.slides, 1):
        p = OUT / f"slide-{i:02d}.png"
        render_slide(slide, p)
        paths.append(p)
        print("wrote", p)
    print("done", len(paths), "slides")


if __name__ == "__main__":
    main()
