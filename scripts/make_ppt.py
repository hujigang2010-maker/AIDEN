"""生成《上海商办楼宇与产业园区市场数据平台》PPT 路演稿（资本/产品视角）。"""
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import config

NAVY = RGBColor(0x1F, 0x4E, 0x79)
BLUE = RGBColor(0x2E, 0x75, 0xB6)
ORANGE = RGBColor(0xE0, 0x7B, 0x39)
GREY = RGBColor(0x59, 0x59, 0x59)
LIGHT = RGBColor(0xF2, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Noto Sans CJK SC"

SW, SH = Inches(13.333), Inches(7.5)
CH = config.CHART_DIR

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def _set_font(run, size, bold=False, color=GREY):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT


def rect(slide, x, y, w, h, color):
    sp = slide.shapes.add_shape(1, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (txt, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = txt
        _set_font(r, size, bold, color)
        p.space_after = Pt(4)
    return tb


def bullets(slide, x, y, w, h, items, size=16, color=GREY, gap=8):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = "•  " + it
        _set_font(r, size, False, color)
        p.space_after = Pt(gap)
    return tb


def header(slide, title, idx):
    rect(slide, 0, 0, SW, Inches(1.1), NAVY)
    rect(slide, 0, Inches(1.1), SW, Inches(0.06), ORANGE)
    textbox(slide, Inches(0.6), Inches(0.18), Inches(11), Inches(0.8),
            [(title, 26, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
    textbox(slide, Inches(12.2), Inches(0.18), Inches(0.9), Inches(0.8),
            [(f"{idx:02d}", 20, True, ORANGE)], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def footer(slide):
    textbox(slide, Inches(0.6), Inches(7.05), Inches(12), Inches(0.35),
            [("易居房地产研究院 × AI 数据获取方 ｜ 上海商办楼宇与产业园区市场数据平台",
              9, False, GREY)])


def slide_blank():
    return prs.slides.add_slide(BLANK)


# ---- 封面 ----
def cover():
    s = slide_blank()
    rect(s, 0, 0, SW, SH, NAVY)
    rect(s, 0, Inches(4.0), SW, Inches(0.08), ORANGE)
    textbox(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.6),
            [("上海商办楼宇与产业园区", 44, True, WHITE),
             ("市场数据平台", 44, True, WHITE)])
    textbox(s, Inches(0.95), Inches(4.2), Inches(11.5), Inches(1.0),
            [("全域覆盖 · 供需双向追踪 · 企业迁徙监测 · 精准招商", 20, False, RGBColor(0xCF, 0xDD, 0xEE))])
    textbox(s, Inches(0.95), Inches(6.3), Inches(11.5), Inches(0.6),
            [("路演稿 ｜ 易居房地产研究院 联合研发  ｜  202X", 14, False, RGBColor(0xAF, 0xC4, 0xDD))])


# ---- 1 市场痛点 ----
def pain():
    s = slide_blank()
    header(s, "市场痛点：传统商办研究的四大短板", 1)
    cards = [
        ("重核心区 · 轻全域", "只盯核心商务区甲级写字楼，乙级/产业园/街镇载体覆盖不足"),
        ("重资产 · 轻需求", "只看租金、空置、供应量，缺企业行业/规模/迁徙/扩缩租研究"),
        ("重静态 · 轻动态", "缺月度企业迁徙追踪与空间行为变化监测"),
        ("重描述 · 轻落地", "难直接转化为招商清单、企业线索、调价与运营动作"),
    ]
    x0, y0, w, h, gx, gy = Inches(0.6), Inches(1.5), Inches(6.0), Inches(2.4), Inches(0.3), Inches(0.3)
    for i, (t, d) in enumerate(cards):
        cx = x0 + (i % 2) * (w + gx)
        cy = y0 + (i // 2) * (h + gy)
        rect(s, cx, cy, w, h, LIGHT)
        rect(s, cx, cy, Inches(0.12), h, ORANGE)
        textbox(s, cx + Inches(0.35), cy + Inches(0.25), w - Inches(0.6), Inches(0.7),
                [(t, 19, True, NAVY)])
        textbox(s, cx + Inches(0.35), cy + Inches(1.05), w - Inches(0.6), Inches(1.2),
                [(d, 14, False, GREY)])
    footer(s)


# ---- 2 差异化 ----
def diff():
    s = slide_blank()
    header(s, "我们的差异化：从「看资产」到「看企业」", 2)
    textbox(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.6),
            [("两大核心竞争力，构建相对五大行的研究壁垒", 18, True, NAVY)])
    left = ["全域覆盖：16 区 → 街道/镇 → 楼宇 → 企业 四级颗粒",
            "全类型载体：超甲/甲/乙级 + 产业园/科创园/孵化器",
            "供需并重：既研究楼宇供给，也追踪企业需求行为",
            "产城融合：嵌入产业集聚、迁徙、补链视角"]
    right = ["企业级迁徙追踪：工商变更 + 招聘地址 + 新闻多源交叉",
             "代理空置率 + 议价系数：还原真实成交价格体系",
             "AI 赋能：LLM 结构化、NLP 产业标签、链路预测",
             "落地导向：招商线索/调价/改造建议可直接执行"]
    rect(s, Inches(0.6), Inches(2.1), Inches(6.0), Inches(4.4), LIGHT)
    rect(s, Inches(6.9), Inches(2.1), Inches(5.8), Inches(4.4), LIGHT)
    textbox(s, Inches(0.85), Inches(2.3), Inches(5.6), Inches(0.5), [("全域 · 供需 · 产城", 18, True, NAVY)])
    textbox(s, Inches(7.15), Inches(2.3), Inches(5.4), Inches(0.5), [("动态 · 真实 · 落地", 18, True, ORANGE)])
    bullets(s, Inches(0.85), Inches(2.9), Inches(5.6), Inches(3.4), left, size=15)
    bullets(s, Inches(7.15), Inches(2.9), Inches(5.4), Inches(3.4), right, size=15)
    footer(s)


# ---- 3 产品矩阵 ----
def product():
    s = slide_blank()
    header(s, "产品矩阵：报告 + 数据库 + 招商工具", 3)
    tiers = [
        ("免费层 · 树品牌", "年度白皮书、季度市场概览", BLUE),
        ("付费层 · 稳收入", "月度报告、区域专题、竞品分析", NAVY),
        ("高端层 · 高毛利", "定制招商名单、企业线索、项目诊断、资产估值", ORANGE),
    ]
    x0, y, w, h, gx = Inches(0.6), Inches(1.6), Inches(3.9), Inches(3.6), Inches(0.35)
    for i, (t, d, c) in enumerate(tiers):
        cx = x0 + i * (w + gx)
        rect(s, cx, y, w, h, c)
        textbox(s, cx + Inches(0.3), y + Inches(0.35), w - Inches(0.6), Inches(0.9),
                [(t, 20, True, WHITE)])
        bullets(s, cx + Inches(0.3), y + Inches(1.4), w - Inches(0.6), Inches(2.0),
                d.split("、"), size=15, color=WHITE, gap=10)
    textbox(s, Inches(0.6), Inches(5.5), Inches(12), Inches(1.2),
            [("本质：不是「一份报告」，而是「上海产业空间数据平台」——数据产品 + 招商工具 + 决策系统",
              17, True, NAVY)])
    footer(s)


# ---- 4 数据壁垒 ----
def data_layers():
    s = slide_blank()
    header(s, "数据壁垒：五层数据架构", 4)
    layers = [
        ("① 空间供给（底盘）", "高德 POI/AOI · 房源平台 · 五大行", BLUE),
        ("② 企业画像（需求核心）", "企查查/天眼查 · 国家企信系统", BLUE),
        ("③ 交易租赁（最难）", "房源平台 · 裁判文书 · 招投标 · 中介访谈", NAVY),
        ("④ 企业行为（差异化）", "工商变更 · 招聘地址 · 融资/新闻", ORANGE),
        ("⑤ 政策产业", "政府公开 · 重点产业规划 · 园区定位", BLUE),
    ]
    y = Inches(1.45)
    for t, d, c in layers:
        rect(s, Inches(0.6), y, Inches(4.6), Inches(0.95), c)
        textbox(s, Inches(0.8), y + Inches(0.13), Inches(4.3), Inches(0.7),
                [(t, 16, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, Inches(5.5), y + Inches(0.13), Inches(7.2), Inches(0.7),
                [(d, 15, False, GREY)], anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(1.08)
    footer(s)


# ---- 5 技术路径 ----
def tech():
    s = slide_blank()
    header(s, "技术路径：采集 → 清洗 → 匹配 → 建模", 5)
    steps = ["多源采集\n高德/企查查/房源", "AI 清洗\nLLM 结构化·NLP 标签",
             "楼宇-企业匹配\n地理编码·空间Join", "指标与模型\n空置/匹配/招商优先级"]
    x0, y, w, h, gx = Inches(0.6), Inches(2.2), Inches(2.85), Inches(1.8), Inches(0.45)
    for i, st in enumerate(steps):
        cx = x0 + i * (w + gx)
        rect(s, cx, y, w, h, NAVY if i % 2 == 0 else BLUE)
        textbox(s, cx + Inches(0.2), y, w - Inches(0.4), h,
                [(st.split("\n")[0], 17, True, WHITE), (st.split("\n")[1], 13, False, RGBColor(0xCF, 0xDD, 0xEE))],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < 3:
            textbox(s, cx + w, y, gx, h, [("→", 26, True, ORANGE)],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, Inches(0.6), Inches(4.5), Inches(12), Inches(2.2), [
        "三大核心模型：供需匹配模型、租金预测模型、招商优先级模型",
        "成交租金「三重旁证」：裁判文书纠纷标的 + 招投标中标价 + 中介访谈 + 议价系数",
        "更新机制：租金/空置/迁徙 月度，供给/企业 季度，全量画像 年度",
    ], size=16)
    footer(s)


# ---- 6 杨浦试点成果 ----
def pilot():
    h = json.loads((config.DATA_DIR / "headline.json").read_text(encoding="utf-8"))
    s = slide_blank()
    header(s, "试点样板：杨浦区街道级跑通成果", 6)
    kpis = [("载体样本", h["载体总数"]), ("可租面积(万㎡)", h["可租面积万㎡"]),
            ("成交租金(元/㎡/天)", h["平均成交租金"]), ("空置率(%)", h["空置率"]),
            ("企业样本", h["入驻企业样本数"]), ("招商线索", h["招商线索数"])]
    x0, y, w, gx = Inches(0.6), Inches(1.4), Inches(1.95), Inches(0.07)
    for i, (k, v) in enumerate(kpis):
        cx = x0 + i * (w + gx)
        rect(s, cx, y, w, Inches(1.15), LIGHT)
        textbox(s, cx, y + Inches(0.12), w, Inches(0.55), [(str(v), 22, True, NAVY)],
                align=PP_ALIGN.CENTER)
        textbox(s, cx, y + Inches(0.72), w, Inches(0.4), [(k, 11, False, GREY)],
                align=PP_ALIGN.CENTER)
    s.shapes.add_picture(str(CH / "03_rent_vacancy_by_plate.png"), Inches(0.6), Inches(2.8), height=Inches(3.6))
    s.shapes.add_picture(str(CH / "05_inout_ratio.png"), Inches(7.0), Inches(2.8), height=Inches(3.6))
    textbox(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.4),
            [("注：当前为演示数据模式产出；填入真实 API Key 后由同一管线一键刷新为实采数据。", 11, False, ORANGE)])
    footer(s)


# ---- 7 商业化 ----
def business():
    s = slide_blank()
    header(s, "商业化路径与定价", 7)
    rows = [
        ("政府招商部门", "定制区域报告 + 招商线索", "50–100 万/年"),
        ("产业园运营方", "月度市场简报 + 竞品分析", "20–50 万/年"),
        ("企业选址部门", "选址评估 + 对标分析", "单次 5–15 万"),
        ("资产持有方/基金", "资产估值 + 市场研判", "按项目 10–30 万"),
        ("公开发布（品牌）", "季度概览版", "获客引流"),
    ]
    y = Inches(1.5)
    rect(s, Inches(0.6), y, Inches(12.1), Inches(0.6), NAVY)
    for j, t in enumerate(["客户类型", "产品形态", "定价参考"]):
        textbox(s, Inches(0.8) + Inches(4.0) * j, y + Inches(0.08), Inches(4.0), Inches(0.5),
                [(t, 15, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.6)
    for i, r in enumerate(rows):
        rect(s, Inches(0.6), y, Inches(12.1), Inches(0.78), LIGHT if i % 2 else WHITE)
        for j, t in enumerate(r):
            textbox(s, Inches(0.8) + Inches(4.0) * j, y + Inches(0.14), Inches(4.0), Inches(0.6),
                    [(t, 14, j == 0, NAVY if j == 0 else GREY)], anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.78)
    footer(s)


# ---- 8 落地节奏与分工 ----
def roadmap():
    s = slide_blank()
    header(s, "落地节奏与协作分工", 8)
    textbox(s, Inches(0.6), Inches(1.35), Inches(12), Inches(0.5), [("阶梯式采集路线", 18, True, NAVY)])
    phases = ["阶段一\n七大中心城区+浦东核心", "阶段二\n浦东全域补全", "阶段三\n16 区全域数据库"]
    x0, y, w, gx = Inches(0.6), Inches(1.9), Inches(3.9), Inches(0.25)
    for i, p in enumerate(phases):
        cx = x0 + i * (w + gx)
        rect(s, cx, y, w, Inches(1.3), [BLUE, NAVY, ORANGE][i])
        textbox(s, cx + Inches(0.2), y, w - Inches(0.4), Inches(1.3),
                [(p.split("\n")[0], 17, True, WHITE), (p.split("\n")[1], 13, False, WHITE)],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(0.6), Inches(3.6), Inches(12), Inches(0.5), [("分工", 18, True, NAVY)])
    rect(s, Inches(0.6), Inches(4.1), Inches(6.0), Inches(2.4), LIGHT)
    rect(s, Inches(6.9), Inches(4.1), Inches(5.8), Inches(2.4), LIGHT)
    textbox(s, Inches(0.85), Inches(4.25), Inches(5.6), Inches(0.5), [("易居研究院", 17, True, NAVY)])
    bullets(s, Inches(0.85), Inches(4.8), Inches(5.6), Inches(1.6),
            ["口径/区域/模型制定", "数据校验与市场研判", "报告与可视化、品牌背书", "成交租金等渠道壁垒"], size=14)
    textbox(s, Inches(7.15), Inches(4.25), Inches(5.4), Inches(0.5), [("AI 数据获取方", 17, True, ORANGE)])
    bullets(s, Inches(7.15), Inches(4.8), Inches(5.4), Inches(1.6),
            ["合规采集/清洗/标准化", "楼宇-企业匹配、迁徙计算", "衍生指标、数据库与 API", "创新验证体系"], size=14)
    footer(s)


# ---- 9 结尾 ----
def closing():
    s = slide_blank()
    rect(s, 0, 0, SW, SH, NAVY)
    rect(s, 0, Inches(4.2), SW, Inches(0.08), ORANGE)
    textbox(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.4),
            [("打穿数据这条线，", 38, True, WHITE),
             ("即上海最强招商数据平台", 38, True, WHITE)])
    textbox(s, Inches(0.95), Inches(4.5), Inches(11.5), Inches(1.2),
            [("诚邀合作：数据共建 · 渠道共享 · 商业共赢", 20, False, RGBColor(0xCF, 0xDD, 0xEE))])
    footer(s)


cover(); pain(); diff(); product(); data_layers(); tech(); pilot(); business(); roadmap(); closing()
OUT = config.ROOT / "下载版本" / "上海商办楼宇与产业园区市场数据平台-PPT路演稿.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print("saved", OUT, "slides=", len(prs.slides._sldIdLst))
