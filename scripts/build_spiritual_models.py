#!/usr/bin/env python3
"""从多人闲聊录音总结中推导可迁移模型，输出 PPTX + DOCX。

源材料：2026-08-08 约 1 小时线下闲聊智能总结
（前世回溯、佛道双修、职场接管、能量链接、身份业务等松散话题）
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.shared import Cm, Pt as DocPt, RGBColor as DocRGB

ROOT = Path("/workspace")
OUT_DIR = ROOT / "outputs"
DOC_DIR = ROOT / "docs"
OUT_DIR.mkdir(parents=True, exist_ok=True)
DOC_DIR.mkdir(parents=True, exist_ok=True)

PPTX_PATH = OUT_DIR / "闲聊推导模型_灵性与职场交叉框架.pptx"
DOCX_PATH = OUT_DIR / "闲聊推导模型白皮书.docx"
MD_PATH = DOC_DIR / "模型推导说明.md"

# —— 视觉：墨青 + 玉色，避开紫/奶油陶土/报纸风 ——
INK = RGBColor(0x14, 0x2A, 0x32)
JADE = RGBColor(0x2F, 0x8A, 0x78)
MIST = RGBColor(0xE7, 0xEF, 0xF1)
FOG = RGBColor(0xD0, 0xDC, 0xE0)
AMBER = RGBColor(0xC9, 0x8B, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x2A, 0x2E)
GRAY = RGBColor(0x5A, 0x6A, 0x70)
SOFT = RGBColor(0xF4, 0xF8, 0xF9)

FONT = "Microsoft YaHei"


def add_rect(slide, x, y, w, h, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_round(slide, x, y, w, h, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    size=18,
    bold=False,
    color=DARK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font=FONT,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=15, color=DARK, bullet_color=JADE):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.2
        r1 = p.add_run()
        r1.text = "●  "
        r1.font.name = FONT
        r1.font.size = Pt(size)
        r1.font.color.rgb = bullet_color
        r2 = p.add_run()
        r2.text = item
        r2.font.name = FONT
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
    return tb


def header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SW, Inches(0.88), INK)
    add_rect(slide, 0, Inches(0.88), SW, Inches(0.05), JADE)
    add_text(slide, Inches(0.5), Inches(0.18), Inches(12.2), Inches(0.5), title, size=24, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.55), Inches(12.2), Inches(0.3), subtitle, size=11, color=FOG)


def footer(slide, page, total=14):
    add_text(
        slide,
        Inches(0.5),
        Inches(7.15),
        Inches(10),
        Inches(0.28),
        "闲聊推导模型 · 源：2026-08-08 多人线下录音总结",
        size=10,
        color=GRAY,
    )
    add_text(
        slide,
        Inches(11.5),
        Inches(7.15),
        Inches(1.4),
        Inches(0.28),
        f"{page}/{total}",
        size=10,
        color=GRAY,
        align=PP_ALIGN.RIGHT,
    )


def build_pptx():
    global SW, SH, BLANK
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    BLANK = prs.slide_layouts[6]
    TOTAL = 14

    # 1 Cover
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, INK)
    add_rect(s, 0, 0, Inches(0.18), SH, JADE)
    add_text(s, Inches(0.9), Inches(1.5), Inches(11), Inches(0.4), "从松散闲聊中蒸馏结构", size=16, color=JADE)
    add_text(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(1.2), "闲聊推导模型", size=48, bold=True, color=WHITE)
    add_text(
        s,
        Inches(0.9),
        Inches(3.4),
        Inches(11.5),
        Inches(0.6),
        "灵性时空 · 能量链接 · 职场接管 · 身份场域",
        size=20,
        color=FOG,
    )
    add_rect(s, Inches(0.9), Inches(4.2), Inches(2.2), Inches(0.06), AMBER)
    add_text(
        s,
        Inches(0.9),
        Inches(4.6),
        Inches(11),
        Inches(0.8),
        "源材料：2026-08-08 12:37–13:39 · 约 6 人 · 多人闲聊智能总结\n方法：把叙事碎片抽象为可迁移的操作模型与隐喻方程",
        size=14,
        color=FOG,
    )
    add_text(s, Inches(0.9), Inches(6.5), Inches(11), Inches(0.4), "INTERNAL · 概念推演，非事实断言", size=12, color=GRAY)

    # 2 Thesis
    s = prs.slides.add_slide(BLANK)
    header(s, "核心命题", "闲聊不是噪声，而是未标注的知识矿")
    cards = [
        ("表象", "话题松散穿越：\n前世故事 / 项目接管 /\n佛道双修 / 装修楼层"),
        ("结构", "反复出现的判断规则：\n存量、配额、闭环、\n三角验证、命名先验"),
        ("产物", "九个可迁移模型 +\n一套「闲聊蒸馏」方法"),
    ]
    for i, (t, body) in enumerate(cards):
        x = Inches(0.6 + i * 4.2)
        add_round(s, x, Inches(1.4), Inches(3.9), Inches(3.6), MIST)
        add_rect(s, x, Inches(1.4), Inches(3.9), Inches(0.08), JADE)
        add_text(s, x + Inches(0.3), Inches(1.7), Inches(3.3), Inches(0.5), t, size=22, bold=True, color=INK)
        add_text(s, x + Inches(0.3), Inches(2.5), Inches(3.3), Inches(2.2), body, size=16, color=DARK)
    add_text(
        s,
        Inches(0.6),
        Inches(5.3),
        Inches(12),
        Inches(1.2),
        "有意思之处：同一场对话里，灵性叙事与职场方法论共享同构逻辑——\n「先看全量信息 → 再按自有脚本推进 → 能量/信心不足时换剧本」。",
        size=15,
        color=DARK,
    )
    footer(s, 2, TOTAL)

    # 3 Map of 9 models
    s = prs.slides.add_slide(BLANK)
    header(s, "九模型总览", "从录音总结中抽出的结构件")
    models = [
        ("M1", "跨世缘分闭环", "成全→分离→重逢→共业"),
        ("M2", "能量存量-链接方程", "ΔE 取决于相对势能"),
        ("M3", "输出配额与破例", "硬上限 3 + 无聊破例"),
        ("M4", "非线性并行自我场", "时空为伪，并行在场"),
        ("M5", "佛道双修互补栈", "护身底座 + 洞见应用层"),
        ("M6", "全量验货接管法", "好坏全收再自推"),
        ("M7", "面相-业力-职业三角", "三信号闭合验证"),
        ("M8", "命名先验空间选择", "先定名，再选楼层"),
        ("M9", "身份锚点三角", "地理×机构×赛道"),
    ]
    for i, (code, title, sub) in enumerate(models):
        row, col = divmod(i, 3)
        x = Inches(0.55 + col * 4.2)
        y = Inches(1.25 + row * 1.85)
        add_round(s, x, y, Inches(4.0), Inches(1.65), MIST)
        add_text(s, x + Inches(0.25), y + Inches(0.25), Inches(1.0), Inches(0.4), code, size=14, bold=True, color=JADE)
        add_text(s, x + Inches(0.25), y + Inches(0.65), Inches(3.5), Inches(0.4), title, size=18, bold=True, color=INK)
        add_text(s, x + Inches(0.25), y + Inches(1.1), Inches(3.5), Inches(0.35), sub, size=12, color=GRAY)
    footer(s, 3, TOTAL)

    # 4 M1
    s = prs.slides.add_slide(BLANK)
    header(s, "M1 · 跨世缘分闭环", "古代故事的叙事骨架 → 关系动力学模型")
    steps = [
        ("① 促成", "促成男生婚配\n成全之后离开"),
        ("② 分离", "各自分散\n缘未断"),
        ("③ 重逢", "家道中落语境下\n男生自外地归来"),
        ("④ 共业", "再结合、多子\n男生开私塾教学"),
    ]
    for i, (t, body) in enumerate(steps):
        x = Inches(0.5 + i * 3.2)
        add_round(s, x, Inches(1.35), Inches(3.0), Inches(2.5), MIST)
        add_text(s, x + Inches(0.2), Inches(1.55), Inches(2.6), Inches(0.45), t, size=20, bold=True, color=JADE)
        add_text(s, x + Inches(0.2), Inches(2.2), Inches(2.6), Inches(1.4), body, size=14, color=DARK)
        if i < 3:
            add_text(s, x + Inches(2.85), Inches(2.3), Inches(0.4), Inches(0.4), "→", size=22, bold=True, color=AMBER)
    add_rect(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.4), INK)
    add_text(s, Inches(0.8), Inches(4.4), Inches(11.8), Inches(0.4), "可迁移推导", size=16, bold=True, color=JADE)
    add_bullets(
        s,
        Inches(0.8),
        Inches(4.95),
        Inches(11.8),
        Inches(1.5),
        [
            "成全他人不等于退出系统：短暂「媒人」角色可能在下一阶段变成「伴侣」。",
            "落魄语境是重逢条件之一——资源下行时，旧缘更容易重新结网。",
            "气质连续性（今生仍高大威猛/书卷气）被当作跨剧本的身份校验码。",
        ],
        size=14,
        color=WHITE,
        bullet_color=AMBER,
    )
    footer(s, 4, TOTAL)

    # 5 M2
    s = prs.slides.add_slide(BLANK)
    header(s, "M2 · 能量存量-链接方程", "网络是否耗能？取决于相对势能，而非链接本身")
    add_rect(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.3), MIST)
    add_text(
        s,
        Inches(0.8),
        Inches(1.55),
        Inches(11.8),
        Inches(0.9),
        "ΔE ≈ g(存量状态 S) · sign(维度差 D)\n若 D > 0（链接更高位）：每次链接 ΔE > 0（补充）｜若 D ≤ 0 且 S 低：ΔE < 0（消耗）",
        size=16,
        color=INK,
    )
    cols = [
        ("低存量 × 平级链接", "容易感到被抽空\n社交/信息流耗能感强", AMBER),
        ("高存量 × 任意链接", "消耗可承受\n或可忽略", JADE),
        ("任意存量 × 高位链接", "每次都被补能\n「向上连接」策略", JADE),
    ]
    for i, (t, body, accent) in enumerate(cols):
        x = Inches(0.5 + i * 4.2)
        add_round(s, x, Inches(2.9), Inches(4.0), Inches(2.6), SOFT)
        add_rect(s, x, Inches(2.9), Inches(0.12), Inches(2.6), accent)
        add_text(s, x + Inches(0.35), Inches(3.15), Inches(3.5), Inches(0.7), t, size=16, bold=True, color=INK)
        add_text(s, x + Inches(0.35), Inches(4.0), Inches(3.5), Inches(1.2), body, size=14, color=DARK)
    add_text(
        s,
        Inches(0.5),
        Inches(5.8),
        Inches(12.3),
        Inches(0.8),
        "职场隐喻：导师/高势能社群 ≈ 高位链接；无效会议/同温层抱怨 ≈ 平级低效链接。",
        size=14,
        color=GRAY,
    )
    footer(s, 5, TOTAL)

    # 6 M3
    s = prs.slides.add_slide(BLANK)
    header(s, "M3 · 灵性输出配额与破例规则", "硬约束是能量；软破例是「剧本无聊」")
    add_round(s, Inches(0.5), Inches(1.35), Inches(6.0), Inches(5.1), MIST)
    add_text(s, Inches(0.8), Inches(1.6), Inches(5.4), Inches(0.4), "默认规则", size=18, bold=True, color=INK)
    add_bullets(
        s,
        Inches(0.8),
        Inches(2.2),
        Inches(5.4),
        Inches(3.8),
        [
            "每人同类输出上限 ≈ 3 次",
            "限制原因：能量不够（资源约束）",
            "第 4 次属于破例",
            "破例理由：旧剧本无聊，要更换场景",
            "新场景仍校验对方气质连续性",
        ],
        size=15,
    )
    add_round(s, Inches(6.8), Inches(1.35), Inches(6.0), Inches(5.1), INK)
    add_text(s, Inches(7.1), Inches(1.6), Inches(5.4), Inches(0.4), "有意思的推导", size=18, bold=True, color=JADE)
    add_bullets(
        s,
        Inches(7.1),
        Inches(2.2),
        Inches(5.4),
        Inches(3.8),
        [
            "配额制 = 稀缺资源的自我保护",
            "「无聊」被赋予合法破例权",
            "更新剧本本身是能量管理",
            "可映射到产品迭代：重复剧本耗能，换叙事续航",
            "创造性厌倦，比硬扛更理性",
        ],
        size=15,
        color=WHITE,
        bullet_color=AMBER,
    )
    footer(s, 6, TOTAL)

    # 7 M4
    s = prs.slides.add_slide(BLANK)
    header(s, "M4 · 非线性并行自我场", "「时间不是线性的，时空是伪的」")
    add_rect(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.5), MIST)
    add_text(
        s,
        Inches(0.8),
        Inches(1.55),
        Inches(11.8),
        Inches(1.1),
        "主张：另一时空中，在场的三个人可能同时在工作。\n时间线不是单轨录像带，而是可并行的场；「此时此刻」只是局部切片。",
        size=16,
        color=INK,
    )
    boxes = [
        ("线性时间观", "过去→现在→未来\n因果单向\n身份唯一"),
        ("并行场观", "多切片同时在场\n因果可回环\n身份可多重出演"),
        ("实用含义", "不必等「一条线走完」\n可同时经营多剧本\n用气质校验连续性"),
    ]
    for i, (t, body) in enumerate(boxes):
        x = Inches(0.5 + i * 4.2)
        add_round(s, x, Inches(3.2), Inches(4.0), Inches(2.8), SOFT)
        add_text(s, x + Inches(0.3), Inches(3.45), Inches(3.4), Inches(0.5), t, size=18, bold=True, color=JADE)
        add_text(s, x + Inches(0.3), Inches(4.2), Inches(3.4), Inches(1.5), body, size=15, color=DARK)
    footer(s, 7, TOTAL)

    # 8 M5
    s = prs.slides.add_slide(BLANK)
    header(s, "M5 · 佛道双修互补栈", "不互相排斥：底座护身，上层洞见")
    add_round(s, Inches(0.5), Inches(1.35), Inches(6.0), Inches(5.1), MIST)
    add_text(s, Inches(0.8), Inches(1.65), Inches(5.4), Inches(0.4), "道 · 护身层（Firmware）", size=18, bold=True, color=INK)
    add_bullets(
        s,
        Inches(0.8),
        Inches(2.3),
        Inches(5.4),
        Inches(3.5),
        [
            "护身功法评价极高",
            "目标：延年益寿、稳住载体",
            "类比：操作系统 / 身体固件",
            "没有载体稳定，上层洞见难跑",
        ],
        size=15,
    )
    add_round(s, Inches(6.8), Inches(1.35), Inches(6.0), Inches(5.1), MIST)
    add_text(s, Inches(7.1), Inches(1.65), Inches(5.4), Inches(0.4), "佛 · 洞见层（Runtime）", size=18, bold=True, color=INK)
    add_bullets(
        s,
        Inches(7.1),
        Inches(2.3),
        Inches(5.4),
        Inches(3.5),
        [
            "金刚经等经典门槛高",
            "需要能量 + 认知基础",
            "普通人「看不懂」= 缺运行环境",
            "类比：高阶应用要足够内存",
        ],
        size=15,
    )
    footer(s, 8, TOTAL)

    # 9 M6
    s = prs.slides.add_slide(BLANK)
    header(s, "M6 · 全量验货接管法", "失败项目多次易手后的反叙事接手")
    flow = [
        ("背景", "多任项目组未成\n团队信心耗尽"),
        ("接手", "要求好坏产品\n全部提交查验"),
        ("推进", "按自有思路\n重写推进路径"),
        ("转向", "完成后换方向\n进入新群体"),
    ]
    for i, (t, body) in enumerate(flow):
        x = Inches(0.5 + i * 3.2)
        add_round(s, x, Inches(1.4), Inches(3.0), Inches(2.4), MIST)
        add_text(s, x + Inches(0.2), Inches(1.65), Inches(2.6), Inches(0.45), f"{i+1}. {t}", size=18, bold=True, color=JADE)
        add_text(s, x + Inches(0.2), Inches(2.35), Inches(2.6), Inches(1.2), body, size=14, color=DARK)
    add_rect(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.4), INK)
    add_text(s, Inches(0.8), Inches(4.4), Inches(11.8), Inches(0.4), "管理学蒸馏", size=16, bold=True, color=JADE)
    add_bullets(
        s,
        Inches(0.8),
        Inches(4.95),
        Inches(11.8),
        Inches(1.5),
        [
            "不要继承失败叙事：先要全量实物/数据，好坏都看。",
            "信心崩溃是信息与解释权失控的结果；全量验货夺回解释权。",
            "完成后主动换赛道：项目成功 ≠ 永久绑定旧群体。",
        ],
        size=14,
        color=WHITE,
        bullet_color=AMBER,
    )
    footer(s, 9, TOTAL)

    # 10 M7
    s = prs.slides.add_slide(BLANK)
    header(s, "M7 · 面相-业力-职业三角验证", "三信号对齐时，匹配度判断上升")
    tri = [
        ("面相信号", "书卷气重\n文静气质"),
        ("业力/前世推断", "适合教育相关\n老师型角色"),
        ("今生职业现实", "工作属性类似老师\n当场被确认匹配"),
    ]
    for i, (t, body) in enumerate(tri):
        x = Inches(0.7 + i * 4.15)
        add_round(s, x, Inches(1.5), Inches(3.9), Inches(2.6), MIST)
        add_text(s, x + Inches(0.3), Inches(1.8), Inches(3.3), Inches(0.5), t, size=18, bold=True, color=INK)
        add_text(s, x + Inches(0.3), Inches(2.5), Inches(3.3), Inches(1.3), body, size=15, color=DARK)
    add_text(
        s,
        Inches(0.7),
        Inches(4.5),
        Inches(12),
        Inches(1.8),
        "模型形式：若 Face ∩ Karma_guess ∩ Job_now ≠ ∅，则「匹配度高」被强化。\n"
        "这是一种启发式诊断，不是科学证明——但它展示了对话中如何用闭环提高确信感。\n"
        "可迁移用法：招聘/合伙人选择时，外表气质、历史叙事、当前行为三线交叉验证。",
        size=15,
        color=DARK,
    )
    footer(s, 10, TOTAL)

    # 11 M8 + M9
    s = prs.slides.add_slide(BLANK)
    header(s, "M8 / M9 · 命名先验与身份锚点", "先有名字与身份，再落物理与业务")
    add_round(s, Inches(0.5), Inches(1.35), Inches(6.0), Inches(5.1), MIST)
    add_text(s, Inches(0.8), Inches(1.6), Inches(5.4), Inches(0.45), "M8 命名先验空间选择", size=18, bold=True, color=INK)
    add_bullets(
        s,
        Inches(0.8),
        Inches(2.25),
        Inches(5.4),
        Inches(3.8),
        [
            "先确定名称「21」",
            "再选择实际 21 层场地",
            "意图/符号先于物理选址",
            "21 层及以上材料统一 10 斤规格，可做顶光",
            "品牌一致性向下渗透到建材",
        ],
        size=14,
    )
    add_round(s, Inches(6.8), Inches(1.35), Inches(6.0), Inches(5.1), MIST)
    add_text(s, Inches(7.1), Inches(1.6), Inches(5.4), Inches(0.45), "M9 身份锚点三角", size=18, bold=True, color=INK)
    add_bullets(
        s,
        Inches(7.1),
        Inches(2.25),
        Inches(5.4),
        Inches(3.8),
        [
            "地理：复旦大学周边专职",
            "机构：杨浦区科技企业联合中心秘书长",
            "赛道：新级增长 / 不动产 / 人工智能",
            "三角稳住后，灵性社交也可落在同一场域",
            "身份不是单点头衔，是可叠加的坐标系",
        ],
        size=14,
    )
    footer(s, 11, TOTAL)

    # 12 Cross-isomorphism
    s = prs.slides.add_slide(BLANK)
    header(s, "跨域同构：灵性脚本 ≈ 职场脚本", "最有意思的发现不是单个模型，而是它们共享骨架")
    rows = [
        ("灵性侧", "职场侧", "共享骨架"),
        ("看前世/面相，找匹配", "全量验货，夺回解释权", "先收集信号，再下判断"),
        ("能量不足最多出三次", "团队信心耗尽后易手", "资源配额决定可持续输出"),
        ("剧本无聊就换场景", "项目完成后换赛道群体", "厌倦是换轨信号，不是失败"),
        ("链接高位则补能", "进入新势能社群", "向上连接改变 ΔE 符号"),
        ("佛道双修不互斥", "增长+不动产+AI 并行", "多栈互补，而非二选一"),
    ]
    for i, (a, b, c) in enumerate(rows):
        y = Inches(1.25 + i * 0.9)
        bg = INK if i == 0 else (MIST if i % 2 else SOFT)
        tc = WHITE if i == 0 else DARK
        add_rect(s, Inches(0.5), y, Inches(4.0), Inches(0.8), bg)
        add_rect(s, Inches(4.55), y, Inches(4.0), Inches(0.8), bg)
        add_rect(s, Inches(8.6), y, Inches(4.2), Inches(0.8), bg)
        add_text(s, Inches(0.65), y + Inches(0.2), Inches(3.7), Inches(0.45), a, size=13, bold=(i == 0), color=tc)
        add_text(s, Inches(4.7), y + Inches(0.2), Inches(3.7), Inches(0.45), b, size=13, bold=(i == 0), color=tc)
        add_text(s, Inches(8.75), y + Inches(0.2), Inches(3.9), Inches(0.45), c, size=13, bold=(i == 0), color=tc)
    footer(s, 12, TOTAL)

    # 13 Distillation method
    s = prs.slides.add_slide(BLANK)
    header(s, "方法：闲聊知识蒸馏五步", "下次拿到松散录音总结时，可复用的抽取流程")
    steps = [
        ("1", "标金句", "找出带判断动词的句子\n（是/不是/取决于/最多）"),
        ("2", "抽规则", "把金句改写成 if-then\n或上限/条件/例外"),
        ("3", "找闭环", "故事是否有成-分-合\n或多信号交叉验证"),
        ("4", "做映射", "把灵性词换成职场/\n产品/关系同构词"),
        ("5", "写模型名", "给规则起可复述的名字\n便于传播与检验"),
    ]
    for i, (n, t, body) in enumerate(steps):
        x = Inches(0.4 + i * 2.55)
        add_round(s, x, Inches(1.4), Inches(2.4), Inches(4.5), MIST)
        add_rect(s, x + Inches(0.85), Inches(1.7), Inches(0.7), Inches(0.7), JADE)
        add_text(
            s,
            x + Inches(0.85),
            Inches(1.7),
            Inches(0.7),
            Inches(0.7),
            n,
            size=20,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(s, x + Inches(0.15), Inches(2.7), Inches(2.1), Inches(0.5), t, size=16, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.15), Inches(3.4), Inches(2.1), Inches(2.0), body, size=12, color=DARK, align=PP_ALIGN.CENTER)
    footer(s, 13, TOTAL)

    # 14 Closing
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, INK)
    add_rect(s, 0, 0, Inches(0.18), SH, JADE)
    add_text(s, Inches(0.9), Inches(1.8), Inches(11.5), Inches(0.5), "收束", size=16, color=JADE)
    add_text(
        s,
        Inches(0.9),
        Inches(2.4),
        Inches(11.5),
        Inches(1.5),
        "闲聊里最可带走的，不是前世情节本身，\n而是一套关于存量、配额、闭环与向上链接的操作语法。",
        size=26,
        bold=True,
        color=WHITE,
    )
    add_text(
        s,
        Inches(0.9),
        Inches(4.4),
        Inches(11.5),
        Inches(1.2),
        "金句回响：\n「能量消耗与否要看本身存量状态，链接高位能量每次都会给你补充能量。」\n「时间不是线性的，时空是伪的。」",
        size=15,
        color=FOG,
    )
    add_text(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.4), "详见同目录白皮书 DOCX / docs 说明", size=12, color=GRAY)

    prs.save(str(PPTX_PATH))
    print(f"Wrote {PPTX_PATH}")


def set_run_font(run, size=11, bold=False, color=None, name="Microsoft YaHei"):
    run.font.size = DocPt(size)
    run.font.bold = bold
    run.font.name = name
    r = run._element
    r.rPr.rFonts.set(qn("w:eastAsia"), name)
    if color is not None:
        run.font.color.rgb = color


def add_heading_cn(doc, text, level=1):
    p = doc.add_heading(text, level=level)
    for run in p.runs:
        set_run_font(run, size=16 if level == 1 else 13, bold=True, color=DocRGB(0x14, 0x2A, 0x32))
    return p


def add_para(doc, text, *, size=11, bold=False, space_after=8):
    p = doc.add_paragraph()
    run = p.add_run(text)
    set_run_font(run, size=size, bold=bold, color=DocRGB(0x1E, 0x2A, 0x2E))
    p.paragraph_format.space_after = DocPt(space_after)
    p.paragraph_format.line_spacing = 1.35
    return p


def build_docx():
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Cm(2.2)
    section.bottom_margin = Cm(2.2)
    section.left_margin = Cm(2.4)
    section.right_margin = Cm(2.4)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = title.add_run("闲聊推导模型白皮书")
    set_run_font(r, size=22, bold=True, color=DocRGB(0x14, 0x2A, 0x32))

    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = sub.add_run("基于 2026-08-08 多人线下录音智能总结的结构蒸馏")
    set_run_font(r, size=11, color=DocRGB(0x5A, 0x6A, 0x70))

    add_para(
        doc,
        "说明：本文是对一场松散闲聊的概念推演，把叙事碎片整理为可迁移模型。"
        "不将其视为历史事实或灵性真理的证明，而视为「人如何在对话中组织意义」的样本分析。",
        size=10,
    )

    add_heading_cn(doc, "一、为什么值得推导")
    add_para(
        doc,
        "这场录音表面上是前世故事、佛道观点、职场经历、楼层装修的跳跃式闲聊。"
        "但反复出现的不是情节，而是判断规则：能量有没有存量、输出有没有上限、"
        "重逢有没有闭环、匹配有没有三角验证、选址有没有命名先验。"
        "把这些规则抽出来，闲聊就从噪声变成了一套操作语法。"
    )

    add_heading_cn(doc, "二、九个模型")

    models = [
        (
            "M1 跨世缘分闭环",
            "促成婚配 → 分散离开 → 家道中落语境重逢 → 再结合并开私塾教学。"
            "推导：成全他人并非退出关系系统；落魄语境可能成为重逢条件；气质连续性被当作跨剧本校验码。",
        ),
        (
            "M2 能量存量-链接方程",
            "网络链接是否耗能，取决于个人存量状态与链接对象的维度差。"
            "高位链接每次补能；低存量面对平级链接更易消耗。"
            "职场映射：向上连接导师/高势能社群，比同温层消耗更划算。",
        ),
        (
            "M3 输出配额与破例规则",
            "默认同类灵性输出最多三次（能量不够）；第四次因「剧本无聊」破例更换。"
            "推导：硬约束是资源，软破例是创造性厌倦。产品与内容工作同样适用——重复剧本耗能，换叙事续航。",
        ),
        (
            "M4 非线性并行自我场",
            "时间非线性，时空为伪；另一时空中在场三人可能同时在工作。"
            "实用含义：不必单线走完才开下一条；可用「气质连续性」在多剧本中保持身份感。",
        ),
        (
            "M5 佛道双修互补栈",
            "道负责护身延寿（底座），佛家经典需要能量与认知基础（应用层）。"
            "普通人看不懂金刚经，被解释为缺少运行环境，而非经典本身无价值。技术隐喻：Firmware + Runtime。",
        ),
        (
            "M6 全量验货接管法",
            "失败项目多次易手、团队信心耗尽后接手：要求好坏产品全部查验，按自有思路推进，完成后换赛道。"
            "核心：不继承失败叙事；用全量信息夺回解释权；成功后允许解绑旧群体。",
        ),
        (
            "M7 面相-业力-职业三角",
            "书卷气面相 → 适合教育的前世/倾向推断 → 当前工作类似老师 → 匹配度被强化。"
            "形式：三信号交集非空则确信上升。可迁移到招聘与合伙人选择。",
        ),
        (
            "M8 命名先验空间选择",
            "先定名「21」，再选 21 层；上层材料规格统一并可做顶光。"
            "意图与符号先于物理选址，品牌一致性向下渗透到建材。",
        ),
        (
            "M9 身份锚点三角",
            "复旦周边专职 × 杨浦科技企业联合中心秘书长 ×（新级增长 / 不动产 / 人工智能）。"
            "身份是坐标系叠加，而非单点头衔；灵性社交也可落在同一场域。",
        ),
    ]
    for title, body in models:
        add_heading_cn(doc, title, level=2)
        add_para(doc, body)

    add_heading_cn(doc, "三、跨域同构（最有意思的一层）")
    add_para(
        doc,
        "灵性叙事与职场叙事在此共享骨架：先收集信号再判断；资源配额决定可持续输出；"
        "厌倦是换轨信号；向上连接改变能量符号；多栈互补优于二选一。"
        "因此，这场闲聊真正可带走的不是「前世情节」，而是一套可在管理、产品、关系中复用的语法。"
    )

    add_heading_cn(doc, "四、闲聊蒸馏五步法")
    for line in [
        "1. 标金句：找出带判断动词的句子（是/不是/取决于/最多）。",
        "2. 抽规则：改写成 if-then，或「上限 / 条件 / 例外」。",
        "3. 找闭环：故事是否有成-分-合，或多信号交叉验证。",
        "4. 做映射：把灵性词换成职场、产品、关系同构词。",
        "5. 写模型名：给规则起可复述的名字，便于传播与检验。",
    ]:
        add_para(doc, line, space_after=4)

    add_heading_cn(doc, "五、可继续玩的方向")
    add_para(
        doc,
        "1）把 M2/M3 做成「社交能量账本」原型：记录链接对象层级与事后能量感。"
        "2）把 M6 写成项目接管检查清单（全量验货表）。"
        "3）把 M1/M7 做成叙事疗法/教练提问卡（闭环与三角验证）。"
        "4）对更多 TicNote/智能总结样本跑同一蒸馏流程，比较哪些模型稳定复现。"
    )

    add_heading_cn(doc, "六、金句存档")
    for q in [
        "「佛道两家各有各的精华，道家的护身功法很好，能够延年益寿。」",
        "「时间不是线性的，时空是伪的。」",
        "「能量消耗与否要看本身存量状态，链接高位能量每次都会给你补充能量。」",
        "「一般每个人做灵性输出最多三次，能量不够，第四次是因为觉得剧本无聊破例更换。」",
    ]:
        add_para(doc, q, size=11)

    note = doc.add_paragraph()
    r = note.add_run("配套交付：outputs/闲聊推导模型_灵性与职场交叉框架.pptx")
    set_run_font(r, size=10, color=DocRGB(0x5A, 0x6A, 0x70))

    doc.save(str(DOCX_PATH))
    print(f"Wrote {DOCX_PATH}")


def build_md():
    md = """# 闲聊推导模型说明

源：2026-08-08 12:37–13:39 多人线下闲聊智能总结。

## 结论（一句话）

松散闲聊里稳定复现的不是情节，而是**存量、配额、闭环、三角验证、命名先验、向上链接**这套操作语法；灵性脚本与职场脚本同构。

## 交付物

| 文件 | 说明 |
|------|------|
| `outputs/闲聊推导模型_灵性与职场交叉框架.pptx` | 14 页模型推演幻灯片 |
| `outputs/闲聊推导模型白皮书.docx` | 文字版白皮书 |
| `scripts/build_spiritual_models.py` | 生成脚本 |

## 九模型速查

1. **跨世缘分闭环** — 成全→分离→重逢→共业  
2. **能量存量-链接方程** — ΔE 取决于存量与维度差  
3. **输出配额与破例** — 上限 3 + 无聊破例  
4. **非线性并行自我场** — 时空为伪，可并行在场  
5. **佛道双修互补栈** — 护身底座 + 洞见应用层  
6. **全量验货接管法** — 好坏全收再自推、完成后换轨  
7. **面相-业力-职业三角** — 三信号闭合验证  
8. **命名先验空间选择** — 先定名再选楼层  
9. **身份锚点三角** — 地理 × 机构 × 赛道  

## 复现

```bash
python3 scripts/build_spiritual_models.py
```
"""
    MD_PATH.write_text(md, encoding="utf-8")
    print(f"Wrote {MD_PATH}")


if __name__ == "__main__":
    build_pptx()
    build_docx()
    build_md()
