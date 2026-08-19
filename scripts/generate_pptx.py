#!/usr/bin/env python3
"""生成伤情与处理决策简报。"""

from pathlib import Path

from content import NOW_EIGHT, TRANSFER_THREE

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

NAVY = RGBColor(0x0B, 0x2F, 0x5B)
NAVY_DARK = RGBColor(0x06, 0x1E, 0x3C)
GOLD = RGBColor(0xC4, 0xA3, 0x5A)
LIGHT = RGBColor(0xF3, 0xF6, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1F, 0x2A, 0x37)
GRAY = RGBColor(0x5C, 0x6B, 0x7A)
GREEN = RGBColor(0x2F, 0x6B, 0x4F)
RED = RGBColor(0xA6, 0x3D, 0x2F)
AMBER = RGBColor(0xB5, 0x7A, 0x2A)
SOFT = RGBColor(0xE8, 0xEE, 0xF5)

FONT = "Microsoft YaHei"
FOOTER_TEXT = "青岛抚顺路和哈尔滨路路口交通事故 · 内部材料 · 底稿 08-17 补充 08-19 · 伤残尚未鉴定"
TOTAL = 10
SW = Inches(13.333)
SH = Inches(7.5)


def add_rect(slide, x, y, w, h, fill, *, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_round(slide, x, y, w, h, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    shp.adjustments[0] = 0.08
    return shp


def set_run(run, text, *, size=16, bold=False, color=DARK, font=FONT):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    from lxml import etree

    rpr = run._r.get_or_add_rPr()
    ea = rpr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
    if ea is None:
        ea = etree.SubElement(rpr, "{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
    ea.set("typeface", font)


def add_text(slide, x, y, w, h, text, *, size=16, bold=False, color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        set_run(p.add_run(), line, size=size, bold=bold, color=color)
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=14, color=DARK, bullet_color=GOLD, spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = spacing
        set_run(p.add_run(), "●  ", size=size, color=bullet_color)
        set_run(p.add_run(), item, size=size, color=color)
    return tb


def header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SW, Inches(0.95), NAVY)
    add_rect(slide, 0, Inches(0.95), SW, Inches(0.05), GOLD)
    add_text(slide, Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.42), title, size=22, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.58), Inches(12.3), Inches(0.3), subtitle, size=12, color=LIGHT)


def footer(slide, page):
    add_text(slide, Inches(0.5), Inches(7.12), Inches(10.8), Inches(0.28), FOOTER_TEXT, size=10, color=GRAY)
    add_text(slide, Inches(11.6), Inches(7.12), Inches(1.2), Inches(0.28), f"{page} / {TOTAL}", size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, title, body, *, title_color=NAVY, accent=GOLD):
    add_round(slide, x, y, w, h, LIGHT)
    add_rect(slide, x, y, Inches(0.08), h, accent)
    add_text(slide, x + Inches(0.22), y + Inches(0.12), w - Inches(0.35), Inches(0.32), title, size=13, bold=True, color=title_color)
    add_text(slide, x + Inches(0.22), y + Inches(0.44), w - Inches(0.35), h - Inches(0.55), body, size=12, color=DARK)


def build_ppt(output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    blank = prs.slide_layouts[6]

    # 1 封面
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, NAVY_DARK)
    add_rect(s, 0, 0, Inches(0.18), SH, GOLD)
    add_text(s, Inches(0.7), Inches(1.6), Inches(12), Inches(0.4), "山东大学齐鲁医院（青岛）云影像  ×  得到大脑通话整理", size=14, color=GOLD)
    add_text(s, Inches(0.7), Inches(2.1), Inches(12), Inches(1.3), "抚顺路和哈尔滨路路口交通事故\n伤情鉴定与伤残情况", size=36, bold=True, color=WHITE)
    add_text(s, Inches(0.7), Inches(4.5), Inches(12), Inches(0.9), "伤者：胡某 · 男 · 64岁\n结论：影像已定性，伤残尚未鉴定，现在不能报残级", size=18, color=LIGHT)
    add_text(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.3), "内部决策简报  ·  2026年8月17日", size=12, color=GRAY)

    # 2 三句话
    s = prs.slides.add_slide(blank)
    header(s, "先看结论", "病例链接打开后，只需要记住这三句")
    card(s, Inches(0.45), Inches(1.25), Inches(4.0), Inches(5.5), "1. 谁受伤", "只有 64 岁男性胡某。\n\n三份 CT 都是他的：\n右小腿、右足、术后右踝。\n\n家属胡继刚是对接人，不是伤者。\n骑手没有本次伤残材料。", accent=GOLD)
    card(s, Inches(4.65), Inches(1.25), Inches(4.0), Inches(5.5), "2. 伤了什么", "本次外伤（急诊 CT）：\n• 右胫骨远端骨折，移位\n• 右腓骨近端骨折，移位\n• 右踝关节半脱位\n• 距骨内缘、外踝撕脱骨折\n\n8月15日已做内外固定，\n复查固定在位、位线可。", accent=RED, title_color=RED)
    card(s, Inches(8.85), Inches(1.25), Inches(4.0), Inches(5.5), "3. 谁伤残", "现在：谁都还没有伤残等级。\n\nCT 是影像诊断，\n不是伤残鉴定书。\n\n评残要等治疗终结后，\n由司法鉴定机构做。\n现在对外禁止说「几级残」。", accent=AMBER, title_color=AMBER)
    footer(s, 2)

    # 3 三份CT
    s = prs.slides.add_slide(blank)
    header(s, "齐鲁医院三份已审核 CT", "检查号可直接用于病历和理赔材料")
    rows = [
        ("08-14 急诊", "10008056847  右小腿", "胫骨远端+腓骨近端骨折\n踝关节半脱位\n距骨内缘、外踝撕脱骨折\n软组织肿胀", RED),
        ("08-14 急诊", "10008056848  右足", "组成骨完整，无新发骨折\n退行性变、跟骨骨刺\n第一跖骨头囊变\n不计入本次评残", GREEN),
        ("08-15 住院", "10008059257  右踝术后", "手足与显微重建外科\n内外固定在位，位线可\n周围软组织仍肿\n证明手术已在三甲完成", NAVY),
    ]
    for i, (when, no, body, color) in enumerate(rows):
        x = Inches(0.45) + Inches(4.2) * i
        add_round(s, x, Inches(1.25), Inches(4.0), Inches(5.5), LIGHT)
        add_rect(s, x, Inches(1.25), Inches(4.0), Inches(0.7), color)
        add_text(s, x + Inches(0.2), Inches(1.38), Inches(3.6), Inches(0.45), when, size=16, bold=True, color=WHITE)
        add_text(s, x + Inches(0.2), Inches(2.15), Inches(3.6), Inches(0.7), no, size=13, bold=True, color=DARK)
        add_text(s, x + Inches(0.2), Inches(2.9), Inches(3.6), Inches(3.5), body, size=15, color=DARK)
    footer(s, 3)

    # 4 改口
    s = prs.slides.add_slide(blank)
    header(s, "通话必须改口", "得到大脑转写不能当诊断")
    items = [
        ("不要再说", "右脚后跟粉碎性骨折", "右足 CT 无跟骨粉碎骨折，只有骨刺。粉碎在胫骨远端和踝关节。"),
        ("不要再说", "只是腓骨下段断了", "报告写腓骨近端骨折，另有外踝撕脱。两处都要记。"),
        ("不要再说", "打了两根或四根钢钉", "术后 CT：内外固定在位。以手术记录为准。"),
        ("不要再说", "8月15日已经出院", "当天住院 CT 在手足外科，仍在齐鲁术后复查。"),
    ]
    for i, (tag, bad, good) in enumerate(items):
        y = Inches(1.2) + Inches(1.4) * i
        add_round(s, Inches(0.45), y, Inches(12.4), Inches(1.28), LIGHT)
        add_text(s, Inches(0.7), y + Inches(0.12), Inches(2.0), Inches(0.3), tag, size=11, bold=True, color=RED)
        add_text(s, Inches(0.7), y + Inches(0.42), Inches(12.0), Inches(0.32), bad, size=16, bold=True, color=DARK)
        add_text(s, Inches(0.7), y + Inches(0.8), Inches(12.0), Inches(0.35), good, size=13, color=GREEN)
    footer(s, 4)

    # 5 伤残规则
    s = prs.slides.add_slide(blank)
    header(s, "伤残怎么处理", "现在只准备，不报价")
    add_bullets(
        s,
        Inches(0.55),
        Inches(1.3),
        Inches(12.2),
        Inches(5.5),
        [
            "评谁：只评伤者胡某。别人没有残。",
            "何时：治疗终结、临床相对稳定后。骨折内固定常见约伤后 3 个月以上，以鉴定机构书面要求为准。",
            "标准：道路交通事故用《人体损伤致残程度分级》。不要套工伤标准，不要用网上计算器报级。",
            "评什么：右踝功能、行走能力。右足骨刺、退变删掉。",
            "现在对外：禁止说已经几级。面谈不锁总包死价。",
            "医学内部判断：损伤较重，日后若踝关节明显受限，具备启动评残的基础；级数看功能，不看急诊片子。",
        ],
        size=16,
        spacing=1.25,
    )
    footer(s, 5)

    # 6 责任
    s = prs.slides.add_slide(blank)
    header(s, "交警视频排除全责", "抚顺路批发市场 · 抚顺路和哈尔滨路路口 · 认定书尚未出具")
    card(s, Inches(0.45), Inches(1.25), Inches(6.1), Inches(2.5), "三轮车（伤者）", "• 货运三轮，可能被认定机动车\n• 仅抚顺路批发市场内部牌，无青岛市号牌\n• 驾驶人有驾驶资格证书；左转向灯失灵\n• 护栏开口借道，先右再左大弧度", accent=AMBER, title_color=AMBER)
    card(s, Inches(6.75), Inches(1.25), Inches(6.1), Inches(2.5), "美团二轮（女骑手）", "• 未保持安全车距\n• 后方视野更好，有条件规避\n• 无交强险，靠平台保险", accent=NAVY)
    add_round(s, Inches(0.45), Inches(4.0), Inches(12.4), Inches(2.7), LIGHT)
    add_text(s, Inches(0.7), Inches(4.15), Inches(12), Inches(0.35), "可能结果只有两种", size=14, bold=True, color=NAVY)
    add_text(s, Inches(0.7), Inches(4.55), Inches(12), Inches(1.9), "① 双方同等责任　　② 三轮次责、外卖主责\n任何一方全责，交警明确排除。法院看视频也不会超出这个范围。\n8月15–16日家属按全责讨论，从今天起停用。", size=16, color=DARK)
    footer(s, 6)

    # 7 程序与钱
    s = prs.slides.add_slide(blank)
    header(s, "程序和钱", "不要用一般程序去赌全责")
    card(s, Inches(0.45), Inches(1.25), Inches(6.1), Inches(5.5), "简易优先", "不审车辆手续和驾照，快出认定书。\n交警建议先协商走简易，避免己方无牌/可能无证被罚。\n\n前提：问清美团平台保险在该责任比例下能否赔医疗费和日后伤残。", accent=GREEN, title_color=GREEN)
    card(s, Inches(6.75), Inches(1.25), Inches(6.1), Inches(5.5), "钱从保险来", "骑手个人偿付能力弱，交警已提醒判了也可能执行不到。\n\n先要：保单/工号、理赔窗口、是否必须认定书。\n刘孝春不替美团补差额；己方份额和车主过错另计。\n责任方截至 8月16日一分未付。", accent=RED, title_color=RED)
    footer(s, 7)

    # 8 转院三点
    s = prs.slides.add_slide(blank)
    header(s, "转院与保险：核心只三点", "这些材料大量重复。不另找三甲。费用按 4 万元核清")
    accents = (GREEN, NAVY, AMBER)
    for i, t in enumerate(TRANSFER_THREE):
        x = Inches(0.45) + Inches(4.2) * i
        card(
            s,
            x,
            Inches(1.25),
            Inches(4.0),
            Inches(5.5),
            f"{t['n']}. {t['what']}",
            t["how"] + "\n\n完成：" + t["done"],
            accent=accents[i],
            title_color=accents[i],
        )
    footer(s, 8)

    # 9 立即八事
    s = prs.slides.add_slide(blank)
    header(s, "立即八事", "律师看完监控再定程序。现在不谈总赔偿数额、不定残级")
    add_bullets(
        s,
        Inches(0.55),
        Inches(1.2),
        Inches(12.2),
        Inches(5.7),
        [f"{x['n']}. {x['what']}" for x in NOW_EIGHT],
        size=14,
        spacing=1.12,
    )
    footer(s, 9)

    # 10 不要说
    s = prs.slides.add_slide(blank)
    header(s, "面谈禁用清单", "说错一句，责任和评残都会被对方抓住")
    bans = [
        "对方全责 / 已经几级伤残",
        "右脚后跟粉碎性骨折",
        "医保能报就是美团险能报",
        "已经是康复期",
        "装病 / 出院再入院绕考核",
        "现在谈总赔偿数额 / 一次性打包了结",
    ]
    for i, text in enumerate(bans):
        y = Inches(1.2) + Inches(0.85) * i
        add_round(s, Inches(0.45), y, Inches(12.4), Inches(0.75), LIGHT)
        add_rect(s, Inches(0.45), y, Inches(0.12), Inches(0.75), RED)
        add_text(s, Inches(0.8), y + Inches(0.18), Inches(11.8), Inches(0.42), text, size=18, bold=True, color=RED)
    footer(s, 10)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


if __name__ == "__main__":
    root = Path(__file__).resolve().parents[1]
    build_ppt(root / "deliverables" / "青岛抚顺路和哈尔滨路路口交通事故_伤情与伤残简报_20260817.pptx")
