"""生成 PPT：学生实习赋能计划 · 可行性论证方案"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pathlib import Path

# 专业蓝绿主题（商务 / 教育场景）
PRIMARY = RGBColor(0x0B, 0x3D, 0x5C)       # 深蓝
ACCENT = RGBColor(0x1A, 0x7A, 0x6D)        # 青绿
LIGHT = RGBColor(0xE8, 0xF3, 0xF1)         # 浅青底
SOFT = RGBColor(0xD0, 0xE8, 0xE4)          # 柔青
GOLD = RGBColor(0xC4, 0xA3, 0x5A)          # 香槟金点缀
DARK = RGBColor(0x1A, 0x2A, 0x33)          # 深色文字
GREY = RGBColor(0x5A, 0x6A, 0x72)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_DEEP = RGBColor(0x08, 0x2E, 0x45)
ORANGE = RGBColor(0xC0, 0x6A, 0x2F)
RED_SOFT = RGBColor(0xB5, 0x4A, 0x4A)


def set_font(run, name="微软雅黑", size=18, bold=False, color=DARK):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    rFonts = rPr.find(qn("a:ea"))
    if rFonts is None:
        rFonts = rPr.makeelement(qn("a:ea"), {"typeface": name})
        rPr.append(rFonts)
    else:
        rFonts.set("typeface", name)


def add_rect(slide, left, top, width, height, fill=PRIMARY, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_rounded(slide, left, top, width, height, fill=LIGHT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def add_text(slide, left, top, width, height, text, size=18, bold=False,
             color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font="微软雅黑"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        text = [text]
    for i, t in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = t
        set_font(run, name=font, size=size, bold=bold, color=color)
    return tb


def add_bullet_list(slide, left, top, width, height, items, size=14,
                    color=DARK, bullet="•", line_spacing=1.25):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = f"{bullet}  {it}"
        set_font(run, size=size, color=color)
    return tb


def slide_header(slide, title, subtitle=None, page_no=None, total=None):
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.42), fill=PRIMARY)
    add_rect(slide, Inches(0), Inches(0.42), Inches(13.333), Inches(0.05), fill=ACCENT)
    add_text(slide, Inches(0.5), Inches(0.55), Inches(10.5), Inches(0.5),
             title, size=22, bold=True, color=PRIMARY, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.0), Inches(12), Inches(0.3),
                 subtitle, size=12, color=GREY, anchor=MSO_ANCHOR.MIDDLE)
    if page_no and total:
        add_text(slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.3),
                 f"{page_no} / {total}", size=10, color=GREY, align=PP_ALIGN.RIGHT)
    add_rect(slide, Inches(0), Inches(7.35), Inches(13.333), Inches(0.15), fill=PRIMARY)


def make_card(slide, left, top, width, height, title, body_items,
              accent=ACCENT, title_size=14, body_size=11):
    add_rounded(slide, left, top, width, height, fill=LIGHT)
    add_rect(slide, left, top, Inches(0.08), height, fill=accent)
    add_text(slide, left + Inches(0.2), top + Inches(0.12),
             width - Inches(0.3), Inches(0.35),
             title, size=title_size, bold=True, color=accent)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.5),
                    width - Inches(0.3), height - Inches(0.55),
                    body_items, size=body_size, color=DARK)


def make_table(slide, left, top, width, height, headers, rows,
               header_color=PRIMARY, font_size=10):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    table = shape.table
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        cell.margin_left = Inches(0.06)
        cell.margin_right = Inches(0.06)
        cell.margin_top = Inches(0.04)
        cell.margin_bottom = Inches(0.04)
        tf = cell.text_frame
        tf.text = ""
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = h
        set_font(r, size=font_size + 1, bold=True, color=WHITE)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 == 0 else LIGHT
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            tf = cell.text_frame
            tf.text = ""
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(val)
            set_font(r, size=font_size, color=DARK)
    return table


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    TOTAL = 16
    page = [0]

    def new_slide():
        page[0] += 1
        return prs.slides.add_slide(blank), page[0]

    # ── 1 封面 ──
    s, _ = new_slide()
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), fill=BG_DEEP)
    add_rect(s, Inches(0), Inches(0), Inches(0.18), Inches(7.5), fill=ACCENT)
    add_rect(s, Inches(0), Inches(5.2), Inches(13.333), Inches(0.04), fill=GOLD)
    add_text(s, Inches(0.9), Inches(1.4), Inches(11.5), Inches(0.45),
             "可行性论证方案  ·  V1.0", size=16, color=SOFT)
    add_text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.9),
             "学生实习赋能计划", size=44, bold=True, color=WHITE)
    add_text(s, Inches(0.9), Inches(2.95), Inches(11.5), Inches(0.55),
             "面向高中生 · 大学生 · 毕业生的可落地实习服务体系",
             size=20, color=SOFT)
    add_text(s, Inches(0.9), Inches(4.0), Inches(11.5), Inches(0.4),
             "目标：辅助就业转化  ·  提升考学与申报竞争力  ·  构建可持续盈利模式",
             size=14, color=GOLD)
    add_text(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.35),
             "商业企业课题实习 × 就业转化  |  定价机制与盈利测算  |  PPT + Excel 双交付",
             size=13, color=WHITE)
    add_text(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.3),
             "讨论稿  |  2026年7月", size=12, color=GREY)

    # ── 2 目录 ──
    s, p = new_slide()
    slide_header(s, "目录", "方案结构一览", p, TOTAL)
    left_items = [
        "一、项目背景与论证目标",
        "二、三类人群需求分析",
        "三、产品与项目规划总览",
        "四、商业实习项目清单",
        "五、战略取舍：取消开证明业务",
    ]
    right_items = [
        "六、收费标准与定价机制",
        "七、盈利模型与财务测算",
        "八、落地实施路径",
        "九、风险与合规要点",
        "十、论证结论与下一步",
    ]
    for i, t in enumerate(left_items):
        y = 1.6 + i * 0.85
        add_rounded(s, Inches(0.7), Inches(y), Inches(5.6), Inches(0.7), fill=LIGHT)
        add_text(s, Inches(0.95), Inches(y + 0.15), Inches(5.2), Inches(0.4),
                 t, size=16, bold=True, color=PRIMARY, anchor=MSO_ANCHOR.MIDDLE)
    for i, t in enumerate(right_items):
        y = 1.6 + i * 0.85
        add_rounded(s, Inches(6.9), Inches(y), Inches(5.6), Inches(0.7), fill=LIGHT)
        add_text(s, Inches(7.15), Inches(y + 0.15), Inches(5.2), Inches(0.4),
                 t, size=16, bold=True, color=PRIMARY, anchor=MSO_ANCHOR.MIDDLE)

    # ── 3 项目背景 ──
    s, p = new_slide()
    slide_header(s, "一、项目背景与论证目标", "为什么现在做、要证明什么", p, TOTAL)
    make_card(s, Inches(0.5), Inches(1.5), Inches(4.0), Inches(5.2),
              "市场痛点", [
                  "优质实习名额稀缺，信息不对称",
                  "海外/名校申请普遍要求实践证明",
                  "校招履历同质化，缺乏项目背书",
                  "毕业生「有学历无经历」转化难",
                  "低价「开证明」业务风险高、利润薄",
              ], accent=ORANGE)
    make_card(s, Inches(4.7), Inches(1.5), Inches(4.0), Inches(5.2),
              "本方案要回答", [
                  "三类人群各自核心诉求是什么",
                  "高价值商业档如何定价与供给",
                  "收费能否覆盖成本并形成利润",
                  "合作企业清单是否具备可触达性",
                  "6–12 个月如何试点验证",
              ], accent=ACCENT)
    make_card(s, Inches(8.9), Inches(1.5), Inches(3.9), Inches(5.2),
              "成功标准（试点期）", [
                  "签约合作企业 ≥ 15 家",
                  "付费学员 ≥ 200 人",
                  "综合毛利率 ≥ 35%",
                  "证明出具合规率 100%",
                  "就业/升学正向反馈 ≥ 60%",
              ], accent=PRIMARY)

    # ── 4 三类人群总览 ──
    s, p = new_slide()
    slide_header(s, "二、三类人群需求分析 · 总览", "同一产品线，差异化价值主张", p, TOTAL)
    headers = ["人群", "核心目标", "关键交付物", "付费意愿", "转化出口"]
    rows = [
        ["高中生", "海外/名校申请加分", "真实企业课题经历与结业证明", "中高（家长决策）", "大学录取 / 竞赛履历"],
        ["大学生", "丰富履历、校招竞争力", "大厂/科技公司项目实习经历", "中（个人+家庭）", "暑期实习转正 / 校招"],
        ["毕业生", "转入正式就业岗位", "过渡实习 + 岗位对接", "中高（急切就业）", "全职 offer"],
    ]
    make_table(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(3.2),
               headers, rows, font_size=12)
    add_text(s, Inches(0.5), Inches(5.1), Inches(12.3), Inches(1.5),
             [
                 "共性需求：可验证的真实课题经历 + 结业证明（作为交付物，不单独售卖）+ 可写入简历/文书的项目成果",
                 "差异策略：高中生走商业项目档提升申请竞争力；大学生偏品牌与技能；毕业生偏就业转化通道",
             ],
             size=14, color=DARK)

    # ── 5 高中生 ──
    s, p = new_slide()
    slide_header(s, "二、需求分析 · 高中生", "申请海外大学 / 升学竞争力", p, TOTAL)
    make_card(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.2),
              "需求画像", [
                  "部分学生申请美英澳加等海外本科",
                  "文书与活动列表需要「真实可核验」实践",
                  "家长关注：安全、时长、证明正规、推荐信质量",
                  "时间窗口：寒暑假 2–8 周为主",
                  "偏好：科技、研究辅助、真实企业远程/混合项目",
                  "痛点：不知如何找到合规渠道，害怕「水经历」",
              ], accent=PRIMARY, body_size=13)
    make_card(s, Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.2),
              "我们提供的解法", [
                  "进入商业项目档：上海 AI / 具身智能真实课题",
                  "结业证明作为项目交付物（不单独开证明售卖）",
                  "可选加购推荐信（仅限项目学员）",
                  "寒暑假 4–6 周项目节奏匹配申请季",
                  "强调可核验过程材料，拒绝「水经历」",
                  "家长沟通重点：真实企业 + 成果物 + 合规",
              ], accent=ACCENT, body_size=13)

    # ── 6 大学生 ──
    s, p = new_slide()
    slide_header(s, "二、需求分析 · 大学生", "履历丰富与校招准备", p, TOTAL)
    make_card(s, Inches(0.5), Inches(1.5), Inches(4.0), Inches(5.2),
              "核心诉求", [
                  "世界 500 强 / 大厂品牌背书",
                  "可写进简历的项目成果",
                  "导师反馈与技能证书",
                  "与专业相关的岗位匹配",
                  "暑期实习转正机会线索",
              ], accent=PRIMARY)
    make_card(s, Inches(4.7), Inches(1.5), Inches(4.0), Inches(5.2),
              "产品匹配", [
                  "商业实习：腾讯/阿里/谷歌等合作通道",
                  "科技/机器人公司实岗或课题实习",
                  "分层套餐：基础证明 / 标准项目 / 精品导师班",
                  "简历辅导与面试模拟（增值）",
                  "校企联合课题（可选）",
              ], accent=ACCENT)
    make_card(s, Inches(8.9), Inches(1.5), Inches(3.9), Inches(5.2),
              "成功指标", [
                  "完成率 ≥ 90%",
                  "简历采纳率 ≥ 70%",
                  "满意度 NPS ≥ 40",
                  "复购/转介绍 ≥ 20%",
                  "优质渠道复用",
              ], accent=ORANGE)

    # ── 7 毕业生 ──
    s, p = new_slide()
    slide_header(s, "二、需求分析 · 毕业生", "从实习过渡到正式就业", p, TOTAL)
    add_bullet_list(s, Inches(0.6), Inches(1.5), Inches(12), Inches(1.2),
                    [
                        "痛点：毕业即失业窗口期、缺乏可证明的工作经验、校招结束后社招门槛更高",
                        "目标：通过 1–3 个月过渡实习，积累岗位经验并进入全职转化通道",
                    ],
                    size=14, color=DARK, bullet="▸")
    headers = ["阶段", "服务内容", "时长", "关键产出"]
    rows = [
        ["入营评估", "职业方向测评 + 岗位匹配", "3–5 天", "个人发展计划（IDP）"],
        ["过渡实习", "真实项目参与 + 周报复盘", "4–12 周", "实习证明 + 作品集"],
        ["就业冲刺", "内推线索 + 面试辅导", "2–4 周", "面试机会 / Offer"],
        ["售后跟踪", "入职 30/90 天回访", "持续", "转化数据与案例"],
    ]
    make_table(s, Inches(0.5), Inches(2.9), Inches(12.3), Inches(3.6),
               headers, rows, font_size=12)

    # ── 8 产品总览 ──
    s, p = new_slide()
    slide_header(s, "三、产品与项目规划总览", "只做高价值：真实企业课题 + 就业转化", p, TOTAL)
    make_card(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.2),
              "A. 商业实习项目（主业）", [
                  "上海 AI / 具身智能 / 科技公司课题",
                  "世界 500 强 / 知名大厂合作通道",
                  "线下或混合制，含导师与考核",
                  "定价：按品牌层级 × 时长 × 辅导强度",
                  "结业证明为交付物，不单独售卖",
                  "适用：大学生为主，优秀高中生可选",
              ], accent=PRIMARY, title_size=15, body_size=13)
    make_card(s, Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.2),
              "B. 就业转化包（高价值出口）", [
                  "毕业生 1–3 个月过渡实习",
                  "内推线索 + 面试辅导 + 复盘",
                  "作品集与结业证明一体交付",
                  "推荐信仅作项目学员加购项",
                  "盈利点：高客单 + 转化口碑",
                  "不做低价开证明获客",
              ], accent=ACCENT, title_size=15, body_size=13)

    # ── 9 商业项目 · 大厂 ──
    s, p = new_slide()
    slide_header(s, "四、商业实习 · 世界500强与知名大厂", "示例合作池（论证用，落地需逐一签约）", p, TOTAL)
    headers = ["层级", "企业示例", "方向", "建议产品形态"]
    rows = [
        ["S 级 · 全球科技", "Google、Microsoft、Amazon、Meta", "产品/数据/工程", "远程课题 + 导师营"],
        ["S 级 · 互联网大厂", "腾讯、阿里巴巴、字节跳动、美团", "产品运营/研发/市场", "暑期营 / 项目制"],
        ["A 级 · 云与 AI", "火山引擎、阿里云、腾讯云、华为云", "云原生/大模型应用", "课题实训营"],
        ["A 级 · 消费电子", "苹果（渠道侧）、小米、OPPO、vivo", "市场/供应链/用户研究", "短期项目实习"],
        ["B 级 · 综合 500 强", "IBM、SAP、西门子、联合利华、宝洁", "咨询/数字化/品牌", "校企联合课题"],
        ["B 级 · 金融科技", "蚂蚁集团、京东科技、平安科技", "风控/运营/数据分析", "项目制远程实习"],
    ]
    make_table(s, Inches(0.4), Inches(1.5), Inches(12.5), Inches(5.0),
               headers, rows, font_size=11)

    # ── 10 商业项目 · 机器人科技 ──
    s, p = new_slide()
    slide_header(s, "四、商业实习 · 机器人及其他科技公司", "硬科技赛道，差异化履历亮点", p, TOTAL)
    headers = ["类别", "企业示例", "适合人群", "岗位/课题方向"]
    rows = [
        ["人形/工业机器人", "优必选、宇树科技、节卡、埃斯顿", "理工大学生", "算法、机械、嵌入式、测试"],
        ["自动驾驶/出行", "小马智行、文远知行、Momenta", "计算机/车辆相关", "感知标注、仿真、运营支持"],
        ["具身智能/AI 硬件", "智元机器人、银河通用、云深处", "硕本交叉背景", "数据采集、应用 demo"],
        ["企业服务 SaaS", "用友、金蝶、销售易、北森", "商科/信息管理", "实施助理、客户成功"],
        ["半导体/芯片生态", "寒武纪、地平线、兆芯（课题侧）", "电子信息", "资料研究、工具链文档"],
        ["生物科技/医疗 AI", "联影智能、推想科技、数坤", "生医/AI 交叉", "标注、文献、产品助理"],
    ]
    make_table(s, Inches(0.4), Inches(1.5), Inches(12.5), Inches(5.0),
               headers, rows, font_size=11)

    # ── 11 战略取舍：取消开证明 ──
    s, p = new_slide()
    slide_header(s, "五、战略取舍：取消「开证明」业务", "费用低、成本与风险高 → 不如聚焦更高价值业务", p, TOTAL)
    make_card(s, Inches(0.5), Inches(1.5), Inches(4.0), Inches(5.2),
              "停售档位", [
                  "199 元：线上实习·证明版",
                  "699 元：证明 + 申报支持",
                  "1,680 元：证明 + 推荐信协助",
                  "以上三档一律取消、不再报价",
                  "不再以「盖章证明」作为获客入口",
              ], accent=RED_SOFT, body_size=13)
    make_card(s, Inches(4.7), Inches(1.5), Inches(4.0), Inches(5.2),
              "为什么取消", [
                  "客单价过低，难以覆盖真实交付",
                  "核验、客服、舆情与合规成本偏高",
                  "易被归类为「水经历」，伤品牌",
                  "占用企业对接与教务带宽",
                  "机会成本：同样精力可售 5,480+",
              ], accent=ORANGE, body_size=13)
    make_card(s, Inches(8.9), Inches(1.5), Inches(3.9), Inches(5.2),
              "改为做什么", [
                  "主推具身/科技项目 5,480",
                  "上海 AI / 头部通道 11,800",
                  "毕业生就业包 9,800",
                  "证明仅随项目结业交付",
                  "推荐信仅项目学员加购",
              ], accent=ACCENT, body_size=13)

    # ── 12 定价 ──
    s, p = new_slide()
    slide_header(s, "六、收费标准与定价机制", "只保留高价值档，保护单价与品牌", p, TOTAL)
    headers = ["产品线", "套餐", "建议定价（元/人）", "包含内容", "目标人群"]
    rows = [
        ["商业标准", "科技/具身智能项目实习", "3,980–6,980（标准 5,480）", "4–6周真实课题+周报+结业证明", "高中优生/大学生"],
        ["商业进阶", "上海 AI / 大厂通道营", "8,800–15,800（标准 11,800）", "品牌项目+导师+答辩+证明", "大学生/优生"],
        ["就业转化", "毕业生过渡实习包", "6,800–12,800（标准 9,800）", "实习+内推辅导+复盘+证明", "毕业生"],
        ["增值加购", "推荐信（不可单卖）", "+1,500", "仅限在营/结业项目学员", "有申请需要者"],
        ["机构合作", "学校/教培团购", "按人 6–8 折", "批量名额+统一课题交付", "B 端渠道"],
    ]
    make_table(s, Inches(0.35), Inches(1.45), Inches(12.6), Inches(4.2),
               headers, rows, font_size=12)
    add_text(s, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.7),
             "明确不做：任何形式的低价「开证明 / 盖章证明 / 申报邮箱单独售卖」产品。",
             size=14, bold=True, color=ORANGE)

    # ── 13 定价机制说明 ──
    s, p = new_slide()
    slide_header(s, "六、定价机制设计原则", "如何实现可持续盈利", p, TOTAL)
    make_card(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(2.4),
              "价格锚点", [
                  "以「企业品牌 × 课题真实度 × 辅导强度」三维定价",
                  "全部收入来自商业/就业高价值档",
                  "团购与早鸟折扣控制在 20% 以内，保护单价",
              ], accent=PRIMARY, body_size=13)
    make_card(s, Inches(6.8), Inches(1.5), Inches(5.9), Inches(2.4),
              "成本结构（示意）", [
                  "企业通道/导师费：25–40%",
                  "运营与客服：10–15%",
                  "获客（渠道/内容）：15–25%",
                  "平台与合规：5–8%，其余为毛利空间",
              ], accent=ACCENT, body_size=13)
    make_card(s, Inches(0.5), Inches(4.2), Inches(6.0), Inches(2.4),
              "增收杠杆", [
                  "标准档 → 进阶通道升级率目标 20–30%",
                  "增值：推荐信/面试模拟仅项目内加购",
                  "B 端：学校/机构年度框架协议（仍走高价值档）",
              ], accent=ORANGE, body_size=13)
    make_card(s, Inches(6.8), Inches(4.2), Inches(5.9), Inches(2.4),
              "支付与交付节奏", [
                  "定金 30% → 企业确认接收 40% → 结业交付 30%",
                  "未达标可补学/延期一次，控制退费率 < 8%",
                  "企业侧按人头或批次结算，锁定毛利",
              ], accent=PRIMARY, body_size=13)

    # ── 14 盈利测算 ──
    s, p = new_slide()
    slide_header(s, "七、盈利模型与财务测算（试点年）", "取消开证明后：人数可少、客单更高 —— 详见 Excel", p, TOTAL)
    headers = ["情景", "年付费人数", "客单价（均）", "营收（万元）", "综合毛利率", "毛利（万元）"]
    rows = [
        ["保守", "100", "7,500", "75", "35%", "26.3"],
        ["基准", "200", "8,240", "164.8", "38%", "62.6"],
        ["乐观", "300", "9,000", "270", "42%", "113.4"],
    ]
    make_table(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.6),
               headers, rows, font_size=13)
    add_bullet_list(s, Inches(0.6), Inches(4.4), Inches(12), Inches(2.4),
                    [
                        "结构假设（基准）：商业标准 50% / 进阶通道 30% / 就业包 20%；均客单约 8,240 元",
                        "盈亏平衡粗算：固定成本约 28–35 万/年；取消低价档后基准情景毛利更厚，更易覆盖固定成本",
                        "关键敏感因子：企业通道成本、获客 CAC、标准→进阶升级率；Excel 提供可调参数表",
                    ],
                    size=14, color=DARK)

    # ── 15 落地路径 ──
    s, p = new_slide()
    slide_header(s, "八、落地实施路径（6 个月试点）", "可执行里程碑", p, TOTAL)
    headers = ["月份", "阶段", "关键动作", "阶段目标"]
    rows = [
        ["M1", "基建", "高价值三档产品封装、合规话术、官网/表单", "MVP 上线（无开证明档）"],
        ["M2", "供给", "签约上海 AI/具身智能 5–8 家，确认导师池", "可售库存就绪"],
        ["M3", "获客", "高校社团/机构渠道试点，标准档开营", "首批 30 人开营"],
        ["M4", "放量", "AI/头部通道营首期，收集案例与口碑", "付费累计 80 人"],
        ["M5", "转化", "毕业生就业包上线，跟踪 offer 转化", "转化案例 ≥ 10"],
        ["M6", "复盘", "财务与 NPS 复盘，决定是否扩区域/扩赛道", "可行性结论定稿"],
    ]
    make_table(s, Inches(0.4), Inches(1.5), Inches(12.5), Inches(5.0),
               headers, rows, font_size=12)

    # ── 16 风险与结论 ──
    s, p = new_slide()
    slide_header(s, "九–十、风险合规与论证结论", "先守住真实课题与品牌，再谈规模", p, TOTAL)
    make_card(s, Inches(0.5), Inches(1.45), Inches(6.0), Inches(3.3),
              "主要风险与应对", [
                  "企业合作不稳定 → 多供应商池+备选课题",
                  "价格战/水项目竞争 → 强调真实课题与成果物",
                  "退费与纠纷 → 分期交付与服务协议",
                  "数据与未成年人合规 → 隐私告知与监护人同意",
                  "品牌方名称使用不当 → 法务审核对外话术",
              ], accent=RED_SOFT, body_size=12)
    make_card(s, Inches(6.8), Inches(1.45), Inches(5.9), Inches(3.3),
              "可行性结论（建议）", [
                  "取消开证明：避免低价高风险业务",
                  "需求真实：三类人群仍可通过高价值档满足",
                  "盈利路径更清晰：客单提升、毛利更厚",
                  "建议：按基准情景启动 6 个月试点验证",
                  "配套：Excel 报价/清单/测算已同步调整",
              ], accent=ACCENT, body_size=12)
    add_rounded(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.5), fill=PRIMARY)
    add_text(s, Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.4),
             "下一步行动", size=16, bold=True, color=GOLD)
    add_text(s, Inches(0.8), Inches(5.65), Inches(11.7), Inches(0.6),
             "① 确认首批合作企业短名单  ② 锁定三档定价与退费政策  ③ 打通证明编号系统  ④ 选择 1–2 个渠道做冷启动",
             size=14, color=WHITE)

    out = Path(__file__).resolve().parents[1] / "exports"
    out.mkdir(parents=True, exist_ok=True)
    path = out / "学生实习赋能计划_可行性论证方案.pptx"
    prs.save(str(path))
    print(f"已生成: {path}")
    return path


if __name__ == "__main__":
    main()
