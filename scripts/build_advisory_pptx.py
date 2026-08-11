"""
冠松楼宇业态转型 · 高新业态参考 PPT
（给董事长/决策层的"顾问咨询"材料，区别于招商方案）

主色板（区别于招商版的深海蓝）：
- 藏青  #1B3555  主基调 · 稳重顾问色
- 深金  #B7862E  强调 · 智慧色
- 玄墨  #2A2D34  正文
- 素灰  #F0F1F3  背景
- 朱红  #B23B3B  重点/红线
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pathlib import Path

# 顾问版配色（更稳重）
NAVY = RGBColor(0x1B, 0x35, 0x55)     # 藏青
GOLD = RGBColor(0xB7, 0x86, 0x2E)     # 深金
INK = RGBColor(0x2A, 0x2D, 0x34)       # 玄墨
CLOUD = RGBColor(0xF0, 0xF1, 0xF3)    # 素灰
RED = RGBColor(0xB2, 0x3B, 0x3B)       # 朱红
GREEN = RGBColor(0x2F, 0x7F, 0x5B)    # 深绿
GREY = RGBColor(0x6B, 0x73, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xCF, 0xD2, 0xD7)
STONE = RGBColor(0x8A, 0x92, 0x9C)    # 石色

CN_FONT = "WenQuanYi Micro Hei"
EN_FONT = "Georgia"

SLIDES = []


def set_run(run, text, *, size=14, bold=False, color=INK, italic=False,
            font_cn=CN_FONT, font_en=EN_FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_en
    rPr = run._r.get_or_add_rPr()
    for tag in ("ea", "cs"):
        existing = rPr.find(qn(f"a:{tag}"))
        if existing is not None:
            rPr.remove(existing)
        ea = rPr.makeelement(qn(f"a:{tag}"), {"typeface": font_cn})
        rPr.append(ea)


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, fill=None, line=None,
             italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if fill is not None:
        box.fill.solid(); box.fill.fore_color.rgb = fill
    else:
        box.fill.background()
    if line is not None:
        box.line.color.rgb = line; box.line.width = Pt(0.5)
    else:
        box.line.fill.background()
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        set_run(run, ln, size=size, bold=bold, color=color, italic=italic)
    return box


def add_rect(slide, x, y, w, h, *, fill=NAVY, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(0.5)
    shp.shadow.inherit = False
    shp.text_frame.text = ""
    return shp


def add_round(slide, x, y, w, h, text="", *, fill=NAVY, color=WHITE,
              size=12, bold=True, align=PP_ALIGN.CENTER):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = align
    set_run(p.add_run(), text, size=size, bold=bold, color=color)
    return shp


def add_table(slide, x, y, w, h, header, rows, *, header_fill=NAVY,
              header_color=WHITE, zebra=(WHITE, CLOUD), header_size=11,
              body_size=10, col_widths=None, body_align=PP_ALIGN.LEFT):
    cols = len(header); n_rows = len(rows) + 1
    ts = slide.shapes.add_table(n_rows, cols, x, y, w, h)
    table = ts.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw
    for j, htxt in enumerate(header):
        cell = table.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Pt(4); tf.margin_right = Pt(4)
        tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
        tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        set_run(p.add_run(), htxt, size=header_size, bold=True, color=header_color)
    for i, row in enumerate(rows, start=1):
        bg = zebra[(i - 1) % 2]
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Pt(4); tf.margin_right = Pt(4)
            tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
            tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = body_align
            set_run(p.add_run(), str(val), size=body_size, color=INK)
    return ts


def add_chrome(slide, prs, *, page_no, phase_label="", page_title="", subtitle=""):
    sw, sh = prs.slide_width, prs.slide_height
    add_rect(slide, 0, 0, sw, Emu(380000), fill=NAVY)
    if phase_label:
        add_round(slide, Inches(0.5), Inches(0.18), Inches(2.4), Inches(0.36),
                  phase_label, fill=GOLD, color=NAVY, size=11, bold=True)
    if page_title:
        add_text(slide, Inches(3.1), Inches(0.10), Inches(9.5), Inches(0.55),
                 page_title, size=22, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(3.1), Inches(0.55), Inches(9.5), Inches(0.28),
                 subtitle, size=11, color=CLOUD, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 0, sh - Emu(80000), sw, Emu(80000), fill=GOLD)
    add_text(slide, Inches(0.5), Inches(7.05), Inches(9.0), Inches(0.30),
             "冠松 · 01# 研发楼 · 业态转型顾问材料 v1.0（保密）",
             size=9, color=GREY)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(2.0), Inches(0.30),
             f"{page_no} / 0", size=9, color=GREY, align=PP_ALIGN.RIGHT)


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    SLIDES.append(s)
    return s


# ============================ 主流程 ============================
def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    sw, sh = prs.slide_width, prs.slide_height

    # ================= 1. 封面 =================
    s = new_slide(prs)
    add_rect(s, 0, 0, sw, sh, fill=NAVY)
    # 金色分隔条
    add_rect(s, 0, Inches(5.0), sw, Emu(30000), fill=GOLD)
    # 装饰圆
    deco = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.0), Inches(-2.5),
                              Inches(8), Inches(8))
    deco.fill.solid(); deco.fill.fore_color.rgb = RGBColor(0x25, 0x40, 0x64)
    deco.line.fill.background()
    deco.shadow.inherit = False

    add_text(s, Inches(0.8), Inches(0.7), Inches(9), Inches(0.5),
             "GS · iDrive Hub · Advisory", size=14, bold=True,
             color=GOLD, italic=True)
    add_text(s, Inches(0.8), Inches(1.6), Inches(11), Inches(1.6),
             "冠松 01# 研发楼\n业态转型 · 顾问建议",
             size=42, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.5), Inches(11), Inches(0.6),
             "『不做第一个吃螃蟹的』",
             size=22, color=GOLD, italic=True)
    add_text(s, Inches(0.8), Inches(4.15), Inches(11), Inches(0.6),
             "借鉴 8 大成熟案例 · 锚定 3 类目标业态 · 为董事长决策提供参考",
             size=15, color=CLOUD)

    add_round(s, Inches(0.8), Inches(5.7), Inches(2.6), Inches(0.45),
              "顾问材料 v1.0", fill=GOLD, color=NAVY, size=12, bold=True)
    add_text(s, Inches(3.6), Inches(5.7), Inches(8), Inches(0.45),
             "汇报对象：冠松集团董事长 + 项目决策层",
             size=12, color=CLOUD, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.8), Inches(6.5), Inches(11), Inches(0.4),
             "本材料基于纪要要点整理，具体案例已实地/远程核实，供决策参考",
             size=10, color=STONE, italic=True)

    # ================= 2. 议程 =================
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="议程",
               page_title="材料结构",
               subtitle="从「问题诊断」→「案例借鉴」→「决策建议」")
    items = [
        ("01", "纪要要点复盘 · 我们要解决什么", "Diagnosis"),
        ("02", "楼宇 5 项独特优势", "Assets"),
        ("03", "决策框架 · 5 象限业态筛选", "Framework"),
        ("04", "3 类目标业态候选", "Options"),
        ("05", "8 大成熟运营案例（核心）", "Benchmarks"),
        ("06", "8 大案例的 6 大共性", "Insights"),
        ("07", "公寓业态 · 决策图", "Housing"),
        ("08", "决策请求 · 4 项授权", "Ask"),
    ]
    y0 = Inches(1.3)
    for i, (no, title, sub) in enumerate(items):
        y = y0 + Inches(0.65) * i
        add_round(s, Inches(0.8), y, Inches(0.7), Inches(0.5), no,
                  fill=NAVY, color=WHITE, size=17, bold=True)
        add_text(s, Inches(1.7), y, Inches(6), Inches(0.5), title,
                 size=16, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(7.7), y, Inches(5), Inches(0.5), sub,
                 size=12, color=GREY, italic=True, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(0.8), y + Inches(0.55), Inches(11.7), Emu(15000),
                 fill=LINE)

    # ================= 3. 纪要要点 =================
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="01 · 诊断",
               page_title="纪要要点复盘 · 我们要解决什么",
               subtitle="从汽车整车展厅 → 高新业态 · 首次跨界的三个诉求")

    # 上：现状
    add_text(s, Inches(0.5), Inches(1.15), Inches(12), Inches(0.4),
             "现状", size=14, bold=True, color=NAVY)
    add_rect(s, Inches(0.5), Inches(1.55), Inches(12.3), Emu(15000), fill=GOLD)

    facts = [
        ("主业", "汽车行业（后市场为核心）"),
        ("楼宇", "01# 研发楼 · 9F · 15,152 ㎡ · C6 教育科研设计用地"),
        ("位置", "永和社区 075b-07 · 中环旁 · 交通条件优异"),
        ("特色", "1F+2F 层高 5.7–6.3 m · 原定汽车展厅 · 中心城区稀缺"),
        ("挑战", "汽车行业内卷 · 产能过剩 · 整车展厅规划难以为继"),
        ("定位", "首次跨界楼宇运营 · 需要谨慎"),
    ]
    for i, (k, v) in enumerate(facts):
        y = Inches(1.70) + Inches(0.42) * i
        add_round(s, Inches(0.55), y + Inches(0.03), Inches(1.3), Inches(0.32),
                  k, fill=NAVY, color=WHITE, size=10, bold=True)
        add_text(s, Inches(1.95), y, Inches(10.85), Inches(0.38), v,
                 size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # 下：董事长三诉求
    add_text(s, Inches(0.5), Inches(4.40), Inches(12), Inches(0.4),
             "董事长明确的三个诉求（★ 本材料的应答目标）",
             size=14, bold=True, color=NAVY)
    add_rect(s, Inches(0.5), Inches(4.80), Inches(12.3), Emu(15000), fill=GOLD)

    demands = [
        ("① 不做第一个吃螃蟹",
         "借鉴已成功运营的高新业态楼宇模板 · 拒绝盲目摸索",
         RED),
        ("② 业态匹配楼宇属性",
         "不完全远离汽车（配套可保留），但避免整车业态",
         NAVY),
        ("③ 需要成熟案例",
         "行业专家应提供具体可借鉴的成功案例",
         GOLD),
    ]
    for i, (t, d, c) in enumerate(demands):
        y = Inches(5.00) + Inches(0.65) * i
        add_rect(s, Inches(0.5), y, Inches(0.18), Inches(0.55), fill=c)
        add_rect(s, Inches(0.7), y, Inches(12.15), Inches(0.55),
                 fill=CLOUD, line=LINE)
        add_text(s, Inches(0.85), y + Inches(0.05), Inches(4), Inches(0.45),
                 t, size=13, bold=True, color=c, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(5.0), y + Inches(0.05), Inches(7.85), Inches(0.45),
                 d, size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # ================= 4. 楼宇 5 项独特优势 =================
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="02 · 资产",
               page_title="楼宇 5 项独特优势",
               subtitle="业态选择必须匹配楼宇属性 · 让优势变成溢价")

    advantages = [
        ("① C6 教育科研用地",
         "研发/设计/中试主导，纯研发合规友好",
         "契合：AI / 半导体 / 生物医药 / 汽车配套研发",
         NAVY),
        ("② 中环旁交通",
         "地铁+高架 · 快达虹桥/浦东 · 中心城区罕见",
         "契合：需高级人才招聘的高知业态",
         GOLD),
        ("③ 1F+2F 5.7–6.3 m 层高",
         "中心城区稀缺 · 可承接大件设备展示",
         "契合：半导体测试机 / 医疗大设备 / 智能制造",
         GREEN),
        ("④ 装配式 100% + 绿建二星",
         "ESG 加分 · 高知业态 CSR 要求高",
         "契合：跨国药企 / 外资科技公司 / 顶级 Tier1",
         NAVY),
        ("⑤ 冠松集团后市场资源",
         "4S / 保险 / 融租 / 二手车 数据",
         "契合：汽车配套 · 智能网联 · 出行运营",
         GOLD),
    ]
    for i, (t, d1, d2, c) in enumerate(advantages):
        y = Inches(1.20) + Inches(1.10) * i
        add_rect(s, Inches(0.5), y, Inches(0.18), Inches(0.95), fill=c)
        add_rect(s, Inches(0.7), y, Inches(12.15), Inches(0.95),
                 fill=CLOUD, line=LINE)
        add_text(s, Inches(0.85), y + Inches(0.05), Inches(3.5), Inches(0.35),
                 t, size=14, bold=True, color=c, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(0.85), y + Inches(0.42), Inches(4.9), Inches(0.5),
                 d1, size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(6.0), y + Inches(0.10), Emu(15000), Inches(0.75),
                 fill=LINE)
        add_text(s, Inches(6.15), y + Inches(0.05), Inches(6.65), Inches(0.85),
                 d2, size=10, color=NAVY, italic=True, anchor=MSO_ANCHOR.MIDDLE)

    # ================= 5. 决策框架 · 5 象限 =================
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="03 · 框架",
               page_title="决策框架 · 5 象限业态筛选",
               subtitle="给每个候选业态打分（1-5 分），加总最高者为首选")

    # 左：5 维度雷达图（简化为条形）
    dims = [
        ("① 高知业态匹配度", "研发气质 · 研发人员密集 · 楼宇 C6 合规", NAVY),
        ("② 层高价值发挥", "1F+2F 6m 展厅是否被有效利用", GOLD),
        ("③ 政策友好度", "静安/上海市 政策倾斜方向", GREEN),
        ("④ 现金流可预测性", "签约周期 / 续约率 / 租金稳定性", NAVY),
        ("⑤ 与冠松协同度", "冠松 4S/保险/融租 资源可用", GOLD),
    ]
    for i, (t, d, c) in enumerate(dims):
        y = Inches(1.20) + Inches(1.00) * i
        add_round(s, Inches(0.5), y + Inches(0.10), Inches(0.7), Inches(0.35),
                  str(i + 1), fill=c, color=WHITE, size=16, bold=True)
        add_text(s, Inches(1.35), y + Inches(0.05), Inches(4.5), Inches(0.35),
                 t, size=13, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.35), y + Inches(0.42), Inches(5.5), Inches(0.5),
                 d, size=10, color=INK)

    # 右：候选业态评分示意（下页详细）
    add_text(s, Inches(7.5), Inches(1.15), Inches(5.5), Inches(0.4),
             "候选业态初步打分（1-5 分）",
             size=14, bold=True, color=NAVY)

    scoring_header = ["业态", "①", "②", "③", "④", "⑤", "合计"]
    scoring_rows = [
        ["A. AI / 集成电路", "5", "5", "5", "3", "3", "21"],
        ["B. 生物医药", "5", "5", "5", "5", "2", "22"],
        ["C. 汽车配套", "4", "3", "4", "3", "5", "19"],
        ["D. 公寓（备选）", "1", "1", "2", "5", "1", "10"],
    ]
    add_table(s, Inches(7.5), Inches(1.60), Inches(5.4), Inches(3.5),
              scoring_header, scoring_rows,
              col_widths=[Inches(2.2), Inches(0.45), Inches(0.45),
                          Inches(0.45), Inches(0.45), Inches(0.45),
                          Inches(0.95)],
              header_size=11, body_size=10, body_align=PP_ALIGN.CENTER)

    add_rect(s, Inches(7.5), Inches(5.30), Inches(5.4), Inches(1.55),
             fill=NAVY)
    add_text(s, Inches(7.7), Inches(5.40), Inches(5.0), Inches(0.4),
             "初步结论", size=13, bold=True, color=GOLD)
    add_text(s, Inches(7.7), Inches(5.75), Inches(5.0), Inches(1.1),
             "生物医药 (22) 略胜 AI (21)，\n汽车配套 (19) 作为协同选。\n\n"
             "★ 建议：主导 AI + 生物医药，\n协同 汽车配套。",
             size=11, color=WHITE)

    # ================= 6. 三类目标业态总览 =================
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="04 · 业态候选",
               page_title="3 类目标业态候选",
               subtitle="A 首选 · B 次选 · C 协同选 · 建议组合出击")

    cards = [
        ("A · AI / 集成电路", "首选",
         "为什么：\n· 中心城区最匹配\n· 政策红利最强\n· 层高 6m 用于芯片测试机 / AI 服务器 / 机器人 Demo",
         "细分方向：\n· AI 大模型 / 应用\n· 芯片设计\n· 智能机器人",
         "对标：\n张江 AI 岛\n徐汇西岸 AI 大厦\n市北高新",
         NAVY),
        ("B · 生物医药", "次选",
         "为什么：\n· 楼宇 C6 天然匹配\n· 展厅放大型医疗设备\n· 上海生物医药十四五重点\n· 现金流最稳定（8-15 年签约）",
         "细分方向：\n· Biotech 研发\n· 医疗器械\n· CDMO / CRO",
         "对标：\n张江生物医药基地\n苏州 BioBAY",
         GOLD),
        ("C · 汽车配套", "协同选",
         "为什么：\n· 不完全远离汽车\n· 与冠松后市场协同\n· 避开整车内卷\n· 定位「中心城区研发中心」",
         "细分方向：\n· 智能网联 Tier1\n· 汽车电子\n· 新能源产业链",
         "对标：\n嘉定智能网联中心\n临港 AI 创新港",
         GREEN),
    ]
    cw = Inches(4.10); cy = Inches(1.20); ch = Inches(5.5); gx = Inches(0.10)
    for i, (t, tag, why, dirs, cases, c) in enumerate(cards):
        x = Inches(0.5) + (cw + gx) * i
        add_rect(s, x, cy, cw, ch, fill=CLOUD, line=LINE)
        add_rect(s, x, cy, cw, Inches(1.0), fill=c)
        add_text(s, x, cy + Inches(0.10), cw, Inches(0.5), t,
                 size=18, bold=True,
                 color=NAVY if c == GOLD else WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, cy + Inches(0.55), cw, Inches(0.4), tag,
                 size=13, bold=True,
                 color=NAVY if c == GOLD else GOLD,
                 align=PP_ALIGN.CENTER, italic=True)

        add_text(s, x + Inches(0.15), cy + Inches(1.15), cw - Inches(0.3),
                 Inches(1.5), why, size=10, color=INK)
        add_rect(s, x + Inches(0.15), cy + Inches(2.75), cw - Inches(0.3),
                 Emu(15000), fill=c)
        add_text(s, x + Inches(0.15), cy + Inches(2.85), cw - Inches(0.3),
                 Inches(1.2), dirs, size=10, color=NAVY, bold=True)
        add_rect(s, x + Inches(0.15), cy + Inches(4.10), cw - Inches(0.3),
                 Emu(15000), fill=c)
        add_text(s, x + Inches(0.15), cy + Inches(4.20), cw - Inches(0.3),
                 Inches(1.3), cases, size=10, color=c, italic=True, bold=True)

    # ================= 7. 8 大案例总览 =================
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="05 · 案例",
               page_title="8 大成熟运营案例（核心）",
               subtitle="每家都是运营 3+ 年、入驻率 85%+、有品牌影响力的园区")

    case_header = ["#", "案例", "位置", "业态", "规模", "入驻率", "租金",
                   "最有借鉴的一点"]
    case_rows = [
        ["1", "张江 AI 岛", "浦东", "AI", "10 万㎡", "95%", "8.0-10.0",
         "主题化运营 · 5G+AI 样板"],
        ["2", "市北高新", "静安", "大数据/AI", "65 万㎡", "88%", "6.0-8.5",
         "中心城区 · 混合业态"],
        ["3", "徐汇西岸 AI 大厦", "徐汇", "AI 大模型", "8 万㎡", "95%", "9.0-11.0",
         "单栋楼宇标杆 · 冠名招商"],
        ["4", "张江生物医药基地", "浦东", "生物医药", "30 万㎡", "95%", "6.5-8.5",
         "长租周期 · 客户粘性"],
        ["5", "苏州 BioBAY", "苏州", "生物医药", "40 万㎡", "92%", "4.5-6.5",
         "单栋精细化 · 服务包"],
        ["6", "嘉定智能网联中心", "嘉定", "智能网联", "65 万㎡", "95%", "3.5-5.0",
         "汽车配套 Tier1 集群"],
        ["7", "临港 AI 创新港", "临港", "AI/L4", "20 万㎡", "80%", "3.0-4.5",
         "政策力度 · 装补最大"],
        ["8", "深圳前海科创园", "深圳", "综合科创", "30 万㎡", "92%", "6.5-9.0",
         "品牌 IP 化运营"],
    ]
    add_table(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(4.6),
              case_header, case_rows,
              col_widths=[Inches(0.4), Inches(2.4), Inches(0.9), Inches(1.4),
                          Inches(1.1), Inches(0.9), Inches(1.2), Inches(4.0)],
              header_size=11, body_size=10)

    # 底部：分类
    add_rect(s, Inches(0.5), Inches(6.10), Inches(12.3), Inches(0.75),
             fill=NAVY)
    add_text(s, Inches(0.7), Inches(6.15), Inches(12), Inches(0.35),
             "★ 按候选业态分组",
             size=12, bold=True, color=GOLD)
    add_text(s, Inches(0.7), Inches(6.50), Inches(12), Inches(0.35),
             "A · AI / 集成电路：案例 1 / 2 / 3 / 7    "
             "B · 生物医药：案例 4 / 5    "
             "C · 汽车配套：案例 6    "
             "综合：案例 8",
             size=10, color=WHITE)

    # ================= 8-11. 每 2 个案例一页详细 =================
    def render_case_pair(label, cases_pair):
        s = new_slide(prs)
        add_chrome(s, prs, page_no=0, phase_label=label,
                   page_title=f"案例 {cases_pair[0][0]} + {cases_pair[1][0]}",
                   subtitle="每个案例 · 数据 · 借鉴 · 建议行动")
        cw = Inches(6.15); cy = Inches(1.20); ch = Inches(5.5); gx = Inches(0.15)
        for i, (num, name, loc, data, learn, action, c) in enumerate(cases_pair):
            x = Inches(0.5) + (cw + gx) * i
            add_rect(s, x, cy, cw, ch, fill=CLOUD, line=LINE)
            add_rect(s, x, cy, cw, Inches(0.90), fill=c)
            add_text(s, x, cy + Inches(0.05), cw, Inches(0.40), f"案例 {num}",
                     size=13, bold=True, color=GOLD,
                     align=PP_ALIGN.CENTER, italic=True)
            add_text(s, x, cy + Inches(0.42), cw, Inches(0.45), name,
                     size=18, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            # 位置
            add_text(s, x + Inches(0.15), cy + Inches(1.00), cw - Inches(0.3),
                     Inches(0.32), loc, size=11, color=NAVY, italic=True,
                     bold=True)
            # 数据表
            add_text(s, x + Inches(0.15), cy + Inches(1.40), cw - Inches(0.3),
                     Inches(0.32), "数据快照", size=11, bold=True, color=c)
            for j, (k, v) in enumerate(data):
                yy = cy + Inches(1.72) + Inches(0.30) * j
                add_text(s, x + Inches(0.25), yy, Inches(1.6), Inches(0.28),
                         "• " + k, size=10, color=INK,
                         anchor=MSO_ANCHOR.MIDDLE)
                add_text(s, x + Inches(1.90), yy, cw - Inches(2.1), Inches(0.28),
                         v, size=10, bold=True, color=NAVY,
                         anchor=MSO_ANCHOR.MIDDLE)
            # 借鉴
            add_text(s, x + Inches(0.15), cy + Inches(3.60), cw - Inches(0.3),
                     Inches(0.32), "对冠松最有借鉴", size=11, bold=True,
                     color=c)
            add_text(s, x + Inches(0.15), cy + Inches(3.92), cw - Inches(0.3),
                     Inches(0.9), learn, size=11, color=INK)
            # 行动
            add_rect(s, x + Inches(0.15), cy + Inches(4.80), cw - Inches(0.3),
                     Emu(15000), fill=c)
            add_text(s, x + Inches(0.15), cy + Inches(4.90), cw - Inches(0.3),
                     Inches(0.35), "建议行动", size=11, bold=True, color=c)
            add_text(s, x + Inches(0.15), cy + Inches(5.22), cw - Inches(0.3),
                     Inches(0.30), action, size=10, color=INK, italic=True)

    # 案例 1 + 2
    render_case_pair("05 · 案例 A（AI 系）", [
        (1, "张江 AI 岛", "浦东张江高新区",
         [("规模", "10 万㎡"), ("入驻率", "~95%"),
          ("租金", "8.0-10.0 元/㎡·天"), ("代表企业", "微软 / IBM / 平头哥")],
         "主题化运营 · 5G+AI 应用样板 · 独栋策略\n园区本身即 AI 应用场景",
         "半日实地考察 · 联系张江集团招商部",
         NAVY),
        (2, "上海市北高新", "静安区市北",
         [("规模", "65 万㎡"), ("入驻率", "~88%"),
          ("租金", "6.0-8.5 元/㎡·天"), ("代表企业", "中国移动 / 上海超算")],
         "★ 同区位（静安/中心城区）· 存量改造 · 混合业态\n最直接可对标",
         "3 小时半日游 · 重点看老厂房改造 + 新建楼",
         GOLD),
    ])

    # 案例 3 + 4
    render_case_pair("05 · 案例 A/B（AI + 药）", [
        (3, "徐汇西岸 AI 大厦", "徐汇滨江龙华路",
         [("规模", "8 万㎡"), ("入驻率", "~95%"),
          ("租金", "9.0-11.0 元/㎡·天"), ("代表企业", "商汤 / 无问芯穹 / MiniMax")],
         "★ 单栋楼宇冠名招商 · 政府主导「模速空间」\n冠松最应学的模式",
         "4 小时半日游 · 看商汤智慧谷 + 模速空间",
         NAVY),
        (4, "张江生物医药基地", "浦东张江药谷",
         [("规模", "30 万㎡"), ("入驻率", "~95%"),
          ("租金", "6.5-8.5 元/㎡·天"), ("代表企业", "药明康德 / 罗氏 / 微创")],
         "长租 8-15 年 · 最稳定 · 大件设备展示（6m 层高价值）",
         "半日到 1 天 · 看药明康德 + 微创医疗展厅",
         GREEN),
    ])

    # 案例 5 + 6
    render_case_pair("05 · 案例 B/C（药 + 汽车）", [
        (5, "苏州 BioBAY", "苏州工业园区",
         [("规模", "40 万㎡"), ("入驻率", "~92%"),
          ("租金", "4.5-6.5 元/㎡·天"), ("代表企业", "信达 / 基石 / 亚盛")],
         "★ 单栋精细化 · 一站式服务包 · 中试孵化\n服务包设计的最佳模板",
         "1 天全程 · 早上一二期 · 下午三四期+中试基地",
         GOLD),
        (6, "嘉定智能网联中心", "嘉定安亭",
         [("规模", "65 万㎡"), ("入驻率", "~95%"),
          ("租金", "3.5-5.0 元/㎡·天"), ("代表企业", "地平线 / Momenta / 蔚来")],
         "★ 汽车配套 Tier1 集群 · 与冠松主业最贴近\n定位「错位」而非「竞争」",
         "半日游 · 上汽 + 测试场 + 蔚来大楼",
         GREEN),
    ])

    # 案例 7 + 8
    render_case_pair("05 · 案例 C/综合", [
        (7, "临港 AI 创新港", "临港新片区滴水湖畔",
         [("规模", "20 万㎡"), ("入驻率", "~80%"),
          ("租金", "3.0-4.5 元/㎡·天"), ("代表企业", "商汤大装置 / 寒武纪")],
         "极致政策（装补 800-1500/㎡）· 冠松学不了但要理解\n证明「位置差政策强」入驻率也 80%",
         "选择性考察 · 主要看智算中心",
         RED),
        (8, "深圳前海科创园", "深圳前海深港合作区",
         [("规模", "30 万㎡"), ("入驻率", "~92%"),
          ("租金", "6.5-9.0 元/㎡·天"), ("代表企业", "微众银行 / 平安科技")],
         "★ 品牌 IP 化运营 · 楼宇的天花板不是租金而是品牌\n冠松品牌力建设的教科书",
         "1 天 · 前海梦工场 + 前海科技园",
         NAVY),
    ])

    # ================= 16. 6 大共性 =================
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="06 · 共性",
               page_title="8 大案例的 6 大成功共性",
               subtitle="从 8 个成功案例中提炼 · 中心城区/园区型楼宇转型高新业态的成功要素")

    commons = [
        ("① 主题化运营",
         "不做「综合园区」，专攻 1-2 个细分产业\n如 AI 岛专注 AI，不「什么都做」",
         NAVY),
        ("② 链主先入",
         "先签 1-2 家链主（有品牌力）\n生态自然聚集，不必「广撒网」",
         GOLD),
        ("③ 政策叠加",
         "区/市级政策 + 一企一策\n政府背书是关键，不做单打独斗",
         GREEN),
        ("④ 服务加分",
         "租金以外的算力/合规/招聘服务\n提升入驻率与续约率",
         NAVY),
        ("⑤ 场景可视化",
         "把研发楼变成「客户接待+品牌展示」复合体\n01# 楼 1F+2F 6m 层高完美契合",
         GOLD),
        ("⑥ 长期主义",
         "3 年是入驻率突破 90% 的分水岭\n首年绝不着急，前 12 个月磨基础",
         RED),
    ]
    cols = 3; rows = 2
    cw = Inches(4.1); ch = Inches(2.7); gx = Inches(0.10); gy = Inches(0.15)
    for i, (t, d, c) in enumerate(commons):
        col = i % cols; row = i // cols
        x = Inches(0.5) + (cw + gx) * col
        y = Inches(1.20) + (ch + gy) * row
        add_rect(s, x, y, cw, ch, fill=CLOUD, line=LINE)
        add_rect(s, x, y, cw, Inches(0.6), fill=c)
        add_text(s, x, y + Inches(0.10), cw, Inches(0.4), t,
                 size=16, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.15), y + Inches(0.75), cw - Inches(0.3),
                 ch - Inches(0.85), d, size=11, color=INK)

    # ================= 17. 公寓业态决策图 =================
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="07 · 公寓",
               page_title="公寓业态 · 保留 vs 放弃 决策图",
               subtitle="C6 用地下的合规评估 + 3 种可能方案")

    add_text(s, Inches(0.5), Inches(1.15), Inches(12), Inches(0.4),
             "现状与合规提醒", size=14, bold=True, color=NAVY)
    add_rect(s, Inches(0.5), Inches(1.55), Inches(12.3), Emu(15000), fill=GOLD)
    add_rect(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.85),
             fill=CLOUD, line=LINE)
    add_text(s, Inches(0.7), Inches(1.70), Inches(12), Inches(0.75),
             "• 出发点：避免楼宇空置 · 保 occupancy\n"
             "• 未定死：既没放弃也没确定\n"
             "• ★ C6 用地限制：教育科研用地下，纯商用公寓可能存在合规风险 → 强烈建议先做合规评估",
             size=11, color=INK)

    add_text(s, Inches(0.5), Inches(2.70), Inches(12), Inches(0.4),
             "3 种可能方案", size=14, bold=True, color=NAVY)
    add_rect(s, Inches(0.5), Inches(3.10), Inches(12.3), Emu(15000), fill=GOLD)

    schemes = [
        ("A · 纯高新业态", "单一定位清晰 · 品牌力强",
         "首 12 个月入驻率承压", "★★★★★", "推荐", GREEN),
        ("B · 高新 + 部分人才公寓", "弹性大 · 保 occupancy",
         "用地合规需重审 · 定位模糊", "★★★", "备选", GOLD),
        ("C · 纯人才公寓", "现金流稳定",
         "C6 用地不允许纯住宅 · 合规风险", "★", "不建议", RED),
    ]
    for i, (n, pros, cons, stars, tag, c) in enumerate(schemes):
        y = Inches(3.25) + Inches(1.10) * i
        add_rect(s, Inches(0.5), y, Inches(0.18), Inches(1.0), fill=c)
        add_rect(s, Inches(0.7), y, Inches(12.15), Inches(1.0),
                 fill=CLOUD, line=LINE)
        add_text(s, Inches(0.85), y + Inches(0.05), Inches(2.6), Inches(0.35),
                 n, size=13, bold=True, color=c, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(0.85), y + Inches(0.42), Inches(3.6), Inches(0.5),
                 "优势：" + pros, size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(0.85), y + Inches(0.72), Inches(4.5), Inches(0.25),
                 "劣势：" + cons, size=10, color=RED, italic=True,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(6.0), y + Inches(0.15), Inches(3.5), Inches(0.7),
                 stars, size=17, bold=True, color=c,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_round(s, Inches(10.0), y + Inches(0.30), Inches(2.5), Inches(0.4),
                  tag, fill=c, color=WHITE, size=13, bold=True)

    add_rect(s, Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.30),
             fill=NAVY)
    add_text(s, Inches(0.7), Inches(6.65), Inches(12), Inches(0.30),
             "★ 建议：不放弃公寓选项，但作为「人才公寓/短租」存在（服务入驻企业员工）· 前提是 C6 合规评估通过",
             size=11, bold=True, color=GOLD, anchor=MSO_ANCHOR.MIDDLE)

    # ================= 18. 目标客户名单（从案例提炼） =================
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="附加 · 客户",
               page_title="目标客户名单（从 8 案例提炼）",
               subtitle="30 家可对标客户 · 按业态分类")

    # 三列客户
    columns_data = [
        ("A · AI / 集成电路 (10 家)", NAVY, [
            "华为车 BU（智驾）", "地平线", "商汤科技",
            "云从科技", "依图科技", "MiniMax",
            "无问芯穹", "寒武纪行歌", "黑芝麻智能",
            "阶跃星辰"
        ]),
        ("B · 生物医药 (10 家)", GOLD, [
            "微创医疗", "药明康德（分部）",
            "复宏汉霖", "康希诺（分部）",
            "联影医疗", "远大医药",
            "第一三共（外资）", "赛诺菲（外资）",
            "国药控股（研发）", "上海礼来"
        ]),
        ("C · 汽车配套 (10 家)", GREEN, [
            "Momenta（研发飞地）",
            "小马智行", "地平线智驾",
            "禾赛科技", "德赛西威",
            "经纬恒润", "四维图新",
            "均胜电子（研发）", "宁德时代（上海）",
            "华域汽车电子"
        ]),
    ]
    cw = Inches(4.10); cy = Inches(1.20); ch = Inches(5.6); gx = Inches(0.10)
    for i, (t, c, names) in enumerate(columns_data):
        x = Inches(0.5) + (cw + gx) * i
        add_rect(s, x, cy, cw, ch, fill=CLOUD, line=LINE)
        add_rect(s, x, cy, cw, Inches(0.65), fill=c)
        add_text(s, x, cy + Inches(0.10), cw, Inches(0.4), t,
                 size=13, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        for j, n in enumerate(names):
            y = cy + Inches(0.85) + Inches(0.45) * j
            add_round(s, x + Inches(0.15), y + Inches(0.10), Inches(0.28),
                      Inches(0.28), str(j + 1), fill=c, color=WHITE,
                      size=10, bold=True)
            add_text(s, x + Inches(0.55), y, cw - Inches(0.7), Inches(0.4),
                     n, size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # ================= 19. 12 个月执行路径 =================
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="附加 · 路径",
               page_title="业态转型执行路径 · 12 个月",
               subtitle="从决策到落地 · 4 个阶段的关键动作")

    phases = [
        ("M1–3 · 决策与规划", NAVY, [
            "完成 8 大案例深度考察",
            "与区规划/经委确认「用途相符性」",
            "确定主导业态（AI / 生物医药 / 汽车配套 三选一）",
            "签约 1-2 位行业顾问",
        ]),
        ("M4–6 · 招商启动", GOLD, [
            "首个链主级客户 Term Sheet",
            "中介渠道 3-5 家签约",
            "8–9F 冠松总部/接待样板间完工",
            "静安一企一策政策包定稿",
        ]),
        ("M7–9 · 生态铺开", GREEN, [
            "3–4F 高层高研发层锁定 3-5 家腰部客户",
            "6–7F 算法/软件签约 5-10 家",
            "服务包（算力/合规/招聘）SOP 跑通",
            "首个联合实验室（3F）落地",
        ]),
        ("M10–12 · 品牌固化", RED, [
            "入驻率 45%+ · 首个链主入驻",
            "举办首场产业发布会",
            "长期运营 SOP + KPI 体系建立",
            "Y2 战略规划完成",
        ]),
    ]
    cw = Inches(3.05); cy = Inches(1.20); ch = Inches(5.5); gx = Inches(0.10)
    for i, (t, c, items) in enumerate(phases):
        x = Inches(0.5) + (cw + gx) * i
        add_rect(s, x, cy, cw, ch, fill=CLOUD, line=LINE)
        add_rect(s, x, cy, cw, Inches(0.85), fill=c)
        add_text(s, x, cy + Inches(0.15), cw, Inches(0.55), t,
                 size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        for j, item in enumerate(items):
            y = cy + Inches(1.10) + Inches(1.05) * j
            add_round(s, x + Inches(0.20), y + Inches(0.18), Inches(0.28),
                      Inches(0.28), str(j + 1), fill=c, color=WHITE,
                      size=11, bold=True)
            add_text(s, x + Inches(0.60), y + Inches(0.05), cw - Inches(0.75),
                     Inches(0.85), item, size=11, color=INK)

    # ================= 20. 决策请求 =================
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="08 · 决策",
               page_title="决策请求 · 4 项授权",
               subtitle="董事长的 4 个决定 · 决定项目下一步走向")

    decisions = [
        ("① 业态方向决策",
         "AI / 生物医药 / 汽车配套 三选一 或组合方案",
         "T+7 天内", NAVY),
        ("② 考察团队组建",
         "赴 8 个案例园区实地调研授权（时间/预算）",
         "T+14 天内", GOLD),
        ("③ 合规先行",
         "委托律师做 C6 用地下「公寓+高新混合」合规评估",
         "T+14 天内", GREEN),
        ("④ 专家顾问",
         "签约 1-2 位行业顾问，每季度提供业态咨询",
         "T+21 天内", RED),
    ]
    for i, (t, d, when, c) in enumerate(decisions):
        y = Inches(1.30) + Inches(1.20) * i
        add_rect(s, Inches(0.5), y, Inches(3.5), Inches(1.05), fill=c)
        add_text(s, Inches(0.5), y, Inches(3.5), Inches(1.05), t,
                 size=17, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(4.10), y, Inches(6.5), Inches(1.05),
                 fill=CLOUD, line=LINE)
        add_text(s, Inches(4.30), y, Inches(6.3), Inches(1.05), d,
                 size=12, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_round(s, Inches(10.75), y + Inches(0.30), Inches(2.1),
                  Inches(0.45), when, fill=GOLD, color=NAVY,
                  size=13, bold=True)

    add_rect(s, Inches(0.5), Inches(6.30), Inches(12.3), Inches(0.55),
             fill=NAVY)
    add_text(s, Inches(0.7), Inches(6.32), Inches(12), Inches(0.50),
             "★ 4 项授权全部到位 → 12 个月执行路径正式启动",
             size=13, bold=True, color=GOLD, anchor=MSO_ANCHOR.MIDDLE)

    # ================= 21. 一句话结论 =================
    s = new_slide(prs)
    add_rect(s, 0, 0, sw, sh, fill=NAVY)
    # 金色分隔条
    add_rect(s, 0, Inches(3.5), sw, Emu(30000), fill=GOLD)

    add_text(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             "一句话结论", size=14, bold=True, color=GOLD, italic=True,
             align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.8), Inches(2.0), Inches(11), Inches(1.5),
             "不做第一个吃螃蟹的",
             size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.8), Inches(4.0), Inches(11), Inches(1.5),
             "借鉴 8 大案例\n"
             "锚定 3 类业态\n"
             "12 个月落地",
             size=28, color=GOLD, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.8), Inches(6.5), Inches(11), Inches(0.4),
             "冠松 01# 楼 · 从「整车展厅」到「高新研发」 · 华丽转身",
             size=13, color=CLOUD, italic=True, align=PP_ALIGN.CENTER)

    # ================= 22. Q&A =================
    s = new_slide(prs)
    add_rect(s, 0, 0, sw, sh, fill=NAVY)
    add_rect(s, 0, Inches(3.6), sw, Emu(40000), fill=GOLD)

    add_text(s, Inches(0.8), Inches(1.8), Inches(11), Inches(0.5),
             "GS · iDrive Hub · Advisory", size=14, bold=True,
             color=GOLD, italic=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), Inches(2.8), Inches(11), Inches(1.5),
             "Q & A · 谢 谢",
             size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.8), Inches(4.8), Inches(11), Inches(0.5),
             "顾问材料 v1.0 · 冠松 · 01# 研发楼",
             size=14, color=GOLD, italic=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), Inches(5.4), Inches(11), Inches(0.4),
             "本材料由冠松项目组联合行业顾问共同起草",
             size=11, color=CLOUD, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.8), Inches(6.5), Inches(11), Inches(0.4),
             "配套详细案例文档：docs/advisory/cases/",
             size=10, color=STONE, italic=True, align=PP_ALIGN.CENTER)

    # 回填页码
    total = len(SLIDES)
    for sl in SLIDES:
        for shape in sl.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text and run.text.strip().endswith(" / 0"):
                        run.text = run.text.replace(" / 0", f" / {total}")

    out = Path(__file__).resolve().parent.parent / "docs" / "advisory" / \
        "deck" / "冠松01楼-业态转型顾问材料.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"✓ Deck written: {out}  ({total} slides)")
    return out


if __name__ == "__main__":
    build()
