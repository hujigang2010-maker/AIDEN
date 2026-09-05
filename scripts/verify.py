# -*- coding: utf-8 -*-
"""校验交付物是否生成、关键金额与页数是否正确。"""

from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

import content as C

ROOT = Path(__file__).resolve().parent.parent / "output"

FILES = {
    "ppt": ROOT / "港大经管上海中心_外滩产业课堂_联合策划方案.pptx",
    "agreement": ROOT / "港大经管上海中心_联合策划服务合作协议_建议稿.docx",
    "onepager": ROOT / "港大经管上海中心_联合策划_一页纸.docx",
    "xlsx": ROOT / "港大经管上海中心_前期费用与交付清单.xlsx",
    "md": ROOT / "港大经管上海中心_联合策划思路一页纸.md",
    "wx": ROOT / "致潘嘉琰_沟通稿.txt",
    "needs": ROOT / "港大经管上海中心_对方需求整理_20260903.md",
}


def fail(msg: str) -> None:
    print(f"FAIL  {msg}")
    sys.exit(1)


def ok(msg: str) -> None:
    print(f"OK    {msg}")


def text_of_pptx(path: Path) -> str:
    prs = Presentation(path)
    chunks = []
    for slide in prs.slides:
        for shp in slide.shapes:
            if shp.has_text_frame:
                chunks.append(shp.text_frame.text)
            if shp.has_table:
                for row in shp.table.rows:
                    for cell in row.cells:
                        chunks.append(cell.text)
    return "\n".join(chunks)


def text_of_docx(path: Path) -> str:
    doc = Document(path)
    parts = [p.text for p in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                parts.append(cell.text)
    return "\n".join(parts)


def main() -> None:
    for key, path in FILES.items():
        min_size = 200 if key in {"wx"} else 1000
        if not path.exists() or path.stat().st_size < min_size:
            fail(f"{key} 缺失或过小：{path}")
        ok(f"{key} 存在（{path.stat().st_size} bytes）")

    prs = Presentation(FILES["ppt"])
    if len(prs.slides) != 13:
        fail(f"PPT 页数应为 13，实际 {len(prs.slides)}")
    ok("PPT 13 页")
    ppt_text = text_of_pptx(FILES["ppt"])
    for needle in (C.PROJECT_NAME, C.THEIR_CONTACT, "88,000", "前期", "潘嘉琰", "出海", "主任"):
        if needle not in ppt_text:
            fail(f"PPT 缺少关键词：{needle}")
    ok("PPT 含出海口径、接口人与前期费用")

    agr = text_of_docx(FILES["agreement"])
    for needle in (C.FEE_AMOUNT_CN, "一次性", "不是招生代理分成", C.THEIR_EMAIL, "第四条", "出海", "主任"):
        if needle not in agr:
            fail(f"协议缺少关键词：{needle}")
    if agr.count("¥") == 0 and "88000" not in agr.replace(",", ""):
        fail("协议未写明金额")
    ok("协议含前期费用核心条款")

    one = text_of_docx(FILES["onepager"])
    if C.ONE_LINER[:12] not in one:
        fail("一页纸未写入一句话思路")
    ok("一页纸含一句话思路")

    wb = load_workbook(FILES["xlsx"], data_only=False)
    need = ["00-封面", "01-前期费用构成", "02-90天交付清单", "03-首场执行", "04-双方分工", "05-付款与开票", "06-对方需求摘录"]
    if wb.sheetnames != need:
        fail(f"Excel 工作表不符：{wb.sheetnames}")
    ws = wb["01-前期费用构成"]
    if ws["C10"].value != "=SUM(C5:C9)":
        fail(f"费用合计公式不正确：{ws['C10'].value}")
    parts = [ws[f"C{r}"].value for r in range(5, 10)]
    if sum(parts) != C.FEE_AMOUNT:
        fail(f"费用分项合计 {sum(parts)} 不等于 {C.FEE_AMOUNT}")
    ok("Excel 七张表且合计公式正确")

    md = FILES["md"].read_text(encoding="utf-8")
    wx = FILES["wx"].read_text(encoding="utf-8")
    if "8.8" not in wx and "88,000" not in wx and "捌万捌仟" not in wx:
        fail("沟通稿未写费用")
    if "出海" not in wx or C.THEIR_CONTACT not in wx:
        fail("沟通稿未写出海口径或对方姓名")
    if "前期" not in md or "出海" not in md:
        fail("Markdown 未写前期费用或出海")
    needs = FILES["needs"].read_text(encoding="utf-8")
    if "主任" not in needs or "四件事" not in needs:
        fail("需求整理不完整")
    ok("Markdown、沟通稿与需求整理完整")
    print("ALL PASSED")


if __name__ == "__main__":
    main()
