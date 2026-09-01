#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""校验三份交付物：30 场、关键口径、页数。"""

from pathlib import Path
import sys

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "scripts"))
import proposal_data as D  # noqa: E402

from pptx import Presentation
from docx import Document
from openpyxl import load_workbook

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "output"


def main():
    ppt = Presentation(OUT / "同浦汇_30场活动与科技企业服务中心筹备_业务承接策划案.pptx")
    doc = Document(OUT / "同浦汇_30场活动与科技企业服务中心筹备_业务承接策划案.docx")
    wb = load_workbook(OUT / "同浦汇_30场活动与科技企业服务中心筹备_执行台账.xlsx")
    errors = []

    if len(ppt.slides) != 16:
        errors.append(f"PPT 页数 {len(ppt.slides)} ≠ 16")
    ppt_text = []
    for s in ppt.slides:
        for sh in s.shapes:
            if sh.has_text_frame:
                ppt_text.append(sh.text_frame.text)
    blob = "\n".join(ppt_text)
    for k in ["同浦汇", "30 场", "科技企业服务中心", "智能建造", "70%", "38%", "杨浦科创集团"]:
        if k not in blob:
            errors.append(f"PPT 缺少：{k}")

    event_table = None
    for t in doc.tables:
        if t.rows[0].cells[0].text.strip() == "编号" and len(t.rows) >= 31:
            event_table = t
            break
    if event_table is None:
        errors.append("Word 未找到 30 场明细表")
    else:
        n = len(event_table.rows) - 1
        if n != 30:
            errors.append(f"Word 活动行 {n} ≠ 30")

    ws = wb["03-30场总表"]
    codes = [ws.cell(r, 2).value for r in range(4, 34)]
    if codes != [e[0] for e in D.EVENTS]:
        errors.append(f"Excel 场次编号不一致：{codes}")
    ws6 = wb["06-商务付款"]
    if ws6["B15"].value != 30:
        errors.append("Excel 年包基数不是 30 万")
    if ws6["E15"].value != "=B15*C15":
        errors.append("Excel 内部分成公式缺失")
    if "09-下一步" not in wb.sheetnames:
        errors.append("缺少确认表")

    if errors:
        print("FAIL")
        for e in errors:
            print(" -", e)
        sys.exit(1)
    print("OK  PPT=16页  Word活动表=30行  Excel场次=30  分成公式存在")


if __name__ == "__main__":
    main()
