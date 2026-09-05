# -*- coding: utf-8 -*-
"""给港大经管上海中心看的汇报 PPT。口吻按当面汇报来写。"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "output" / "港大经管上海中心_联合合作汇报.pptx"

GREEN = RGBColor(0x00, 0x3D, 0x2E)
GREEN2 = RGBColor(0x0A, 0x5C, 0x46)
GOLD = RGBColor(0xC4, 0xA3, 0x5A)
CREAM = RGBColor(0xF7, 0xF4, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x24, 0x20)
GREY = RGBColor(0x5B, 0x6B, 0x64)
FONT = "微软雅黑"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL = 11


def set_font(run, size, color=DARK, bold=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = rpr.makeelement(qn("a:ea"), {})
        rpr.append(ea)
    ea.set("typeface", FONT)


def add_rect(slide, x, y, w, h, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_round(slide, x, y, w, h, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = 0.08
    except Exception:
        pass
    return shp


def add_text(slide, x, y, w, h, text, size, color=DARK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_font(r, size, color, bold)
    return tb


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def header(slide, title, idx):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.02), GREEN)
    add_rect(slide, 0, Inches(1.02), SLIDE_W, Pt(3), GOLD)
    add_text(slide, Inches(0.55), Inches(0.16), Inches(10.7), Inches(0.7), title, 24, WHITE, True, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(11.5), Inches(0.16), Inches(1.4), Inches(0.7), f"{idx}/{TOTAL}", 12, GOLD, True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def footer(slide):
    add_text(
        slide,
        Inches(0.55),
        Inches(7.12),
        Inches(12.2),
        Inches(0.28),
        "复旦大学住房政策研究中心  ·  上海市杨浦区科技企业联合会  ·  致香港大学经管学院上海中心",
        10,
        GREY,
    )


def bullets(slide, x, y, w, h, items, size=16, gap=10):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.line_spacing = 1.2
        rm = p.add_run()
        rm.text = "·  "
        set_font(rm, size, GOLD, True)
        r = p.add_run()
        r.text = it
        set_font(r, size, DARK)
    return tb


def slide_cover(prs):
    s = blank(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, GREEN)
    add_rect(s, 0, 0, Inches(0.16), SLIDE_H, GOLD)
    add_text(s, Inches(0.85), Inches(1.25), Inches(11.5), Inches(0.4), "致　香港大学经管学院上海中心", 16, GOLD, False)
    add_text(s, Inches(0.85), Inches(1.85), Inches(11.5), Inches(0.9), "联合合作汇报", 44, WHITE, True)
    add_text(
        s,
        Inches(0.85),
        Inches(2.85),
        Inches(11.5),
        Inches(0.55),
        "EMBA招生、外滩出海专题课、请中心主任到场　合在一次合作里",
        20,
        GOLD,
        False,
    )
    add_rect(s, Inches(0.85), Inches(3.55), Inches(2.0), Pt(3), GOLD)
    add_text(
        s,
        Inches(0.85),
        Inches(3.85),
        Inches(11.2),
        Inches(1.2),
        "上次在中心见面之后，把可以马上做的事写成一整包：招生是主项，出海课是这场招生的现场，请主任到场含在同一笔费用里。并附上协议，请中心审阅。",
        18,
        WHITE,
        False,
    )
    add_text(
        s,
        Inches(0.85),
        Inches(5.4),
        Inches(11.2),
        Inches(0.9),
        "复旦大学住房政策研究中心\n上海市杨浦区科技企业联合会",
        16,
        RGBColor(0xD5, 0xE4, 0xDC),
        False,
    )
    add_text(s, Inches(0.85), Inches(6.55), Inches(11.2), Inches(0.35), "2026年9月　　联系人：潘嘉琰老师", 14, GOLD, False)


def slide_why(prs):
    s = blank(prs)
    header(s, "这份材料是什么", 2)
    add_round(s, Inches(0.5), Inches(1.28), Inches(12.3), Inches(1.55), CREAM)
    add_text(
        s,
        Inches(0.75),
        Inches(1.4),
        Inches(11.8),
        Inches(1.3),
        "潘老师：上次谈到华东、华中的招生和市场，主事是 EMBA；也谈到出海，以及请主任到现场看一看。回来后写成一整包：招生是主项，出海课是现场，请主任到场含在同一笔八万八千元里。",
        17,
        DARK,
        False,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    cards = [
        ("第一件　招生", "按港大 EMBA 常见门槛初筛不少于八十人，请三十到四十人到外滩。课后由招生老师单独沟通。录取仍由学院决定，不按人头分成。"),
        ("第二件　出海课", "这场课的主题放在企业出海和赴港上市，同时就是招生现场。不另办一场空的说明会，也不另收费。"),
        ("第三件　请主任", "我们办的出海活动请主任看一次。这一项含在同一笔八万八千元里，不另收费、不单开报价。"),
    ]
    for i, (t, d) in enumerate(cards):
        y = Inches(3.02 + i * 1.28)
        add_round(s, Inches(0.5), y, Inches(12.3), Inches(1.16), WHITE)
        add_rect(s, Inches(0.5), y, Inches(2.15), Inches(1.16), GREEN if i < 2 else GREEN2)
        add_text(s, Inches(0.5), y, Inches(2.15), Inches(1.16), t, 15, WHITE, True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(2.85), y + Inches(0.12), Inches(9.7), Inches(0.92), d, 15, DARK, False, anchor=MSO_ANCHOR.MIDDLE)
    footer(s)


def slide_fit(prs):
    s = blank(prs)
    header(s, "我们为什么来谈这件事", 3)
    bullets(
        s,
        Inches(0.6),
        Inches(1.4),
        Inches(12.1),
        Inches(5.3),
        [
            "潘老师负责华东、华中 EMBA 的招生和市场。我们长期在上海办企业和校友活动，手里有一批企业负责人。",
            "中心在上海已经开始运转。要把合适的人请到外滩来，听老师讲，再由招生老师单独谈。",
            "港大在企业出海、赴港上市这件事上，内地学校很难替代。所以这场招生的主题放在出海上，不办空的说明会。",
            "若以后还要在学院层面合作，需要主任先看过现场。请主任出席我们办的出海活动，和招生、出海课写在同一包里，不另收费。",
        ],
        size=17,
        gap=14,
    )
    footer(s)


def slide_emba(prs):
    s = blank(prs)
    header(s, "EMBA 招生这一项", 4)
    rows = [
        ("名单", "协议生效、费用到账后七日内提交初稿，不少于八十人。双方看过再定，名单共同使用。"),
        ("门槛", "企业负责人或具备相应支付能力的人选，按港大 EMBA 常见要求初筛。具体条件以中心告知为准。"),
        ("到场", "从名单里请三十到四十人到外滩中心上课。不保证每一个人都到，但到场的人按上述门槛来。"),
        ("课后", "由潘老师或中心指定的招生老师单独沟通。我们在七日内把到场情况和有意向的人整理好交给中心。"),
        ("录取", "面试、录取、学位仍由学院决定。我们不承诺录取人数，也不按人头或学费分成。"),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(1.28 + i * 1.08)
        add_round(s, Inches(0.5), y, Inches(12.3), Inches(0.96), WHITE if i % 2 == 0 else CREAM)
        add_text(s, Inches(0.7), y, Inches(1.5), Inches(0.96), k, 16, GREEN, True, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(2.4), y, Inches(10.1), Inches(0.96), v, 15, DARK, False, anchor=MSO_ANCHOR.MIDDLE)
    footer(s)


def slide_class(prs):
    s = blank(prs)
    header(s, "外滩出海专题课", 5)
    rows = [
        ("时间", "协议签订后一个月左右，选一个工作日下午，四点左右结束。"),
        ("地点", "香港大学经管学院上海中心，外滩 SOHO F 栋。"),
        ("人数", "三十到四十人。同时就是上面说的招生现场，不另办说明会。"),
        ("来宾", "正在考虑出海或赴港上市、且符合 EMBA 初筛的企业负责人。"),
        ("内容", "港大老师或校友讲出海和港股路径；企业之间交换看法；课后由招生老师单独沟通。"),
        ("分工", "中心提供场地、老师、招生安排。我们负责请人、现场，会后把名单整理好交给中心。"),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(1.28 + i * 0.88)
        add_round(s, Inches(0.5), y, Inches(12.3), Inches(0.78), WHITE if i % 2 == 0 else CREAM)
        add_text(s, Inches(0.7), y, Inches(1.5), Inches(0.78), k, 16, GREEN, True, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(2.4), y, Inches(10.1), Inches(0.78), v, 15, DARK, False, anchor=MSO_ANCHOR.MIDDLE)
    footer(s)


def slide_director(prs):
    s = blank(prs)
    header(s, "请中心主任到场", 6)
    add_round(s, Inches(0.5), Inches(1.32), Inches(12.3), Inches(2.05), CREAM)
    add_text(
        s,
        Inches(0.8),
        Inches(1.48),
        Inches(11.7),
        Inches(1.75),
        "这一项不是单独向中心报价。八万八千元买的是前面的 EMBA 招生和外滩出海课；请主任到场含在里面，不另收费。\n\n我们近期有出海主题活动，会请总领事或企业负责人到场。希望请潘老师代为邀请主任出席一次，直接看到我们请到的人和现场。来不来，以中心的安排为准。",
        17,
        DARK,
        False,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    bullets(
        s,
        Inches(0.7),
        Inches(3.6),
        Inches(12.0),
        Inches(3.2),
        [
            "我们发出书面邀请，写明时间、地点和到场嘉宾。",
            "主任因公务不能来，不影响这一项算已完成，也不因此减费用或另收费。",
            "港大本部出海课程怎么推广、原创谷怎么合作、学生在上海实习参访，本期不写进交付，以后另说。",
        ],
        size=17,
        gap=12,
    )
    footer(s)


def slide_90(prs):
    s = blank(prs)
    header(s, "九十天怎么排", 7)
    steps = [
        ("第 1 周", "签订协议，支付费用；定下外滩课日期；中心告知招生门槛，以及课后由哪位老师接。"),
        ("第 2–3 周", "按门槛拿出拟邀请名单，不少于八十人，双方看过再定。"),
        ("第 4–7 周", "办外滩出海专题课，课后招生老师单独沟通；同时发出请主任出席的邀请。"),
        ("第 8–12 周", "把课上情况和有意向的人交给中心，由招生老师继续跟。后面若要续办或谈别的事，再另写。"),
    ]
    for i, (t, d) in enumerate(steps):
        y = Inches(1.35 + i * 1.3)
        add_round(s, Inches(0.5), y, Inches(12.3), Inches(1.15), WHITE)
        add_rect(s, Inches(0.5), y, Inches(2.4), Inches(1.15), GREEN)
        add_text(s, Inches(0.5), y, Inches(2.4), Inches(1.15), t, 18, WHITE, True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(3.15), y, Inches(9.4), Inches(1.15), d, 17, DARK, False, anchor=MSO_ANCHOR.MIDDLE)
    footer(s)


def slide_roles(prs):
    s = blank(prs)
    header(s, "各自做什么", 8)
    add_round(s, Inches(0.45), Inches(1.35), Inches(6.1), Inches(5.4), WHITE)
    add_rect(s, Inches(0.45), Inches(1.35), Inches(6.1), Inches(0.7), GREEN)
    add_text(s, Inches(0.45), Inches(1.35), Inches(6.1), Inches(0.7), "香港大学经管学院上海中心", 18, WHITE, True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    bullets(
        s,
        Inches(0.7),
        Inches(2.25),
        Inches(5.6),
        Inches(4.2),
        [
            "提供外滩中心场地和名义",
            "安排老师或校友",
            "告知 EMBA 初筛条件，课后招生跟进",
            "面试和录取仍由学院决定",
            "请潘老师协助邀请中心主任",
        ],
        size=16,
        gap=10,
    )
    add_round(s, Inches(6.8), Inches(1.35), Inches(6.1), Inches(5.4), WHITE)
    add_rect(s, Inches(6.8), Inches(1.35), Inches(6.1), Inches(0.7), GREEN2)
    add_text(s, Inches(6.8), Inches(1.35), Inches(6.1), Inches(0.7), "住房政策研究中心、杨浦科企联", 18, WHITE, True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    bullets(
        s,
        Inches(7.05),
        Inches(2.25),
        Inches(5.6),
        Inches(4.2),
        [
            "按门槛筛名单、请人、现场",
            "会后把纪要和有意向的人交给中心",
            "安排请主任出席的那一场活动",
            "开票、收款",
        ],
        size=16,
        gap=10,
    )
    footer(s)


def slide_fee(prs):
    s = blank(prs)
    header(s, "费用：三件事一笔", 9)
    add_round(s, Inches(0.5), Inches(1.28), Inches(12.3), Inches(1.48), GREEN)
    add_text(s, Inches(0.8), Inches(1.36), Inches(11.7), Inches(0.32), "前期工作费用（一整包）", 14, GOLD, True)
    add_text(
        s,
        Inches(0.8),
        Inches(1.68),
        Inches(11.7),
        Inches(0.9),
        "人民币捌万捌仟元整（88,000 元）　　协议生效后十个工作日内一次付清，开增值税发票",
        18,
        WHITE,
        True,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(s, Inches(0.6), Inches(2.95), Inches(6.0), Inches(0.36), "这八万八千元包括", 16, GREEN, True)
    bullets(
        s,
        Inches(0.55),
        Inches(3.35),
        Inches(6.0),
        Inches(3.45),
        [
            "九十天筹备，以及按 EMBA 门槛组织名单",
            "拟邀请名单不少于八十人",
            "外滩出海专题课一场（同时作招生现场）",
            "课后纪要和有意向的人",
            "请主任出席一次出海活动（含在包内，不另收费）",
        ],
        size=15,
        gap=6,
    )
    add_text(s, Inches(6.9), Inches(2.95), Inches(6.0), Inches(0.36), "不包括", 16, GREEN, True)
    bullets(
        s,
        Inches(6.85),
        Inches(3.35),
        Inches(6.0),
        Inches(3.45),
        [
            "学费、报名费，也不按人头分成",
            "第二场及以后（如需续办，每场六万八千元，另签）",
            "港大本部出海课程的招生",
            "原创谷的日常事务",
            "就请主任到场再单开一笔费用",
        ],
        size=15,
        gap=6,
    )
    footer(s)


def slide_ask(prs):
    s = blank(prs)
    header(s, "请中心定几件事", 10)
    items = [
        "外滩专题课放在哪一天；课后招生老师由谁接、怎么接。",
        "主任方便出席的大致时间。来不了也没有关系，我们按书面邀请办理。",
        "协议由哪一家主体签，发票抬头是谁，款什么时候付。",
    ]
    for i, t in enumerate(items):
        y = Inches(1.4 + i * 1.35)
        add_round(s, Inches(0.55), y, Inches(12.2), Inches(1.2), WHITE)
        add_round(s, Inches(0.8), y + Inches(0.28), Inches(0.65), Inches(0.65), GOLD)
        add_text(s, Inches(0.8), y + Inches(0.28), Inches(0.65), Inches(0.65), str(i + 1), 20, GREEN, True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.7), y + Inches(0.2), Inches(10.7), Inches(0.8), t, 18, DARK, False, anchor=MSO_ANCHOR.MIDDLE)
    add_text(
        s,
        Inches(0.6),
        Inches(5.6),
        Inches(12.1),
        Inches(1.2),
        "协议文本随这份汇报一并送上。签好、费用到账后七日内，我们提交按 EMBA 门槛初筛的拟邀请名单初稿。\n联系人：潘嘉琰　　(86) 180 1860 6086　　jyanpan@hku.hk",
        15,
        GREY,
        False,
        align=PP_ALIGN.CENTER,
    )
    footer(s)


def slide_close(prs):
    s = blank(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, GREEN)
    add_rect(s, 0, 0, Inches(0.16), SLIDE_H, GOLD)
    add_text(s, Inches(0.9), Inches(2.05), Inches(11.5), Inches(0.7), "请中心审阅。有要改的地方，我们按中心的意见改。", 24, WHITE, True)
    add_text(
        s,
        Inches(0.9),
        Inches(2.9),
        Inches(11.5),
        Inches(1.1),
        "招生、出海课、请主任到场，合在一次合作里，费用一笔。",
        20,
        GOLD,
        False,
    )
    add_text(
        s,
        Inches(0.9),
        Inches(4.15),
        Inches(11.5),
        Inches(1.2),
        "复旦大学住房政策研究中心\n上海市杨浦区科技企业联合会\n2026年9月",
        18,
        GOLD,
        False,
    )
    add_text(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(0.8), "随附：联合策划服务合作协议", 16, WHITE, False)


def build(path: Path | None = None) -> Path:
    path = path or OUT
    path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide_cover(prs)
    slide_why(prs)
    slide_fit(prs)
    slide_emba(prs)
    slide_class(prs)
    slide_director(prs)
    slide_90(prs)
    slide_roles(prs)
    slide_fee(prs)
    slide_ask(prs)
    slide_close(prs)
    prs.save(path)
    return path


if __name__ == "__main__":
    out = Path(sys.argv[1]) if len(sys.argv) > 1 else OUT
    print(f"已生成 {build(out)}")
