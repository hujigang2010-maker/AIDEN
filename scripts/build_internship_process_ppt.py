"""生成 PPT：实习推送流程与收费说明（上海 AI / 具身智能企业名录）

针对反馈单独成册：
1. 企业名录（以上海人工智能 + 具身智能为主）
2. 明确收费项目（学生端）
3. 推送至企业的成本测算
4. 端到端清晰流程
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pathlib import Path

PRIMARY = RGBColor(0x0B, 0x3D, 0x5C)
ACCENT = RGBColor(0x1A, 0x7A, 0x6D)
LIGHT = RGBColor(0xE8, 0xF3, 0xF1)
SOFT = RGBColor(0xD0, 0xE8, 0xE4)
GOLD = RGBColor(0xC4, 0xA3, 0x5A)
DARK = RGBColor(0x1A, 0x2A, 0x33)
GREY = RGBColor(0x5A, 0x6A, 0x72)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_DEEP = RGBColor(0x08, 0x2E, 0x45)
ORANGE = RGBColor(0xC0, 0x6A, 0x2F)
RED_SOFT = RGBColor(0xB5, 0x4A, 0x4A)


def set_font(run, name="微软雅黑", size=18, bold=False, color=DARK):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    rFonts = rPr.find(qn("a:ea"))
    if rFonts is None:
        rFonts = rPr.makeelement(qn("a:ea"), {"typeface": name})
        rPr.append(rFonts)
    else:
        rFonts.set("typeface", name)


def add_rect(slide, left, top, width, height, fill=PRIMARY, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_rounded(slide, left, top, width, height, fill=LIGHT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def add_text(slide, left, top, width, height, text, size=18, bold=False,
             color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        text = [text]
    for i, t in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = t
        set_font(run, size=size, bold=bold, color=color)
    return tb


def add_bullet_list(slide, left, top, width, height, items, size=14,
                    color=DARK, bullet="•", line_spacing=1.2):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = f"{bullet}  {it}"
        set_font(run, size=size, color=color)
    return tb


def slide_header(slide, title, subtitle=None, page_no=None, total=None):
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.42), fill=PRIMARY)
    add_rect(slide, Inches(0), Inches(0.42), Inches(13.333), Inches(0.05), fill=ACCENT)
    add_text(slide, Inches(0.5), Inches(0.55), Inches(11), Inches(0.5),
             title, size=22, bold=True, color=PRIMARY, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.0), Inches(12), Inches(0.3),
                 subtitle, size=12, color=GREY, anchor=MSO_ANCHOR.MIDDLE)
    if page_no and total:
        add_text(slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.3),
                 f"{page_no} / {total}", size=10, color=GREY, align=PP_ALIGN.RIGHT)
    add_rect(slide, Inches(0), Inches(7.35), Inches(13.333), Inches(0.15), fill=PRIMARY)


def make_card(slide, left, top, width, height, title, body_items,
              accent=ACCENT, title_size=14, body_size=11):
    add_rounded(slide, left, top, width, height, fill=LIGHT)
    add_rect(slide, left, top, Inches(0.08), height, fill=accent)
    add_text(slide, left + Inches(0.2), top + Inches(0.12),
             width - Inches(0.3), Inches(0.35),
             title, size=title_size, bold=True, color=accent)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.5),
                    width - Inches(0.3), height - Inches(0.55),
                    body_items, size=body_size, color=DARK)


def make_table(slide, left, top, width, height, headers, rows,
               header_color=PRIMARY, font_size=10):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    table = shape.table
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        cell.margin_left = Inches(0.05)
        cell.margin_right = Inches(0.05)
        cell.margin_top = Inches(0.03)
        cell.margin_bottom = Inches(0.03)
        tf = cell.text_frame
        tf.text = ""
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = h
        set_font(r, size=font_size + 1, bold=True, color=WHITE)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 == 0 else LIGHT
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            tf = cell.text_frame
            tf.text = ""
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(val)
            set_font(r, size=font_size, color=DARK)
    return table


def flow_box(slide, left, top, width, height, num, title, lines, fill=PRIMARY):
    add_rounded(slide, left, top, width, height, fill=fill)
    add_text(slide, left + Inches(0.1), top + Inches(0.08), width - Inches(0.2), Inches(0.28),
             f"{num}  {title}", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.1), top + Inches(0.4), width - Inches(0.2), height - Inches(0.45),
             lines, size=10, color=SOFT, align=PP_ALIGN.CENTER)


def arrow_right(slide, left, top):
    add_text(slide, left, top, Inches(0.35), Inches(0.35),
             "→", size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    TOTAL = 12
    page = [0]

    def new_slide():
        page[0] += 1
        return prs.slides.add_slide(blank), page[0]

    # ── 1 封面 ──
    s, _ = new_slide()
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), fill=BG_DEEP)
    add_rect(s, Inches(0), Inches(0), Inches(0.18), Inches(7.5), fill=ACCENT)
    add_rect(s, Inches(0), Inches(5.15), Inches(13.333), Inches(0.04), fill=GOLD)
    add_text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.4),
             "专项说明稿  ·  配套可行性论证方案", size=15, color=SOFT)
    add_text(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(0.85),
             "企业名录 · 收费项目 · 推送流程", size=40, bold=True, color=WHITE)
    add_text(s, Inches(0.9), Inches(3.1), Inches(11.5), Inches(0.5),
             "以上海人工智能企业 + 具身智能企业为主", size=20, color=SOFT)
    add_text(s, Inches(0.9), Inches(4.0), Inches(11.5), Inches(0.4),
             "回答三个问题：有哪些企业可推？学生收多少？推过去成本多少、怎么走？",
             size=15, color=GOLD)
    add_text(s, Inches(0.9), Inches(5.5), Inches(11.5), Inches(0.35),
             "学生实习赋能计划  |  讨论稿 V1.0  |  2026年7月", size=13, color=WHITE)

    # ── 2 目录 / 问题回应 ──
    s, p = new_slide()
    slide_header(s, "本册要讲清楚的三件事", "对应现场反馈：企业名单、收费、流程与推送成本", p, TOTAL)
    cards = [
        ("① 企业名录", PRIMARY, [
            "上海 AI 企业优先池",
            "具身智能 / 机器人企业池",
            "对接状态可更新",
            "收生后按匹配表推送",
        ]),
        ("② 收费项目", ACCENT, [
            "学生端收什么、收多少",
            "分公益 / 商业 / 就业三档",
            "含与不含边界写清",
            "支付节点与退费原则",
        ]),
        ("③ 推送成本与流程", ORANGE, [
            "每个学生推到企业花多少",
            "成本构成一目了然",
            "从报名到出证明七步走",
            "岗位对接人与 SLA",
        ]),
    ]
    for i, (title, color, items) in enumerate(cards):
        x = 0.6 + i * 4.2
        make_card(s, Inches(x), Inches(1.6), Inches(3.9), Inches(5.0),
                  title, items, accent=color, title_size=18, body_size=14)

    # ── 3 企业池总览 ──
    s, p = new_slide()
    slide_header(s, "一、企业名录总览", "已有供给侧基础：以上海本地 AI / 具身智能为主优先对接", p, TOTAL)
    headers = ["类别", "示例企业（上海及本地生态）", "可推岗位/课题", "建议优先"]
    rows = [
        ["大模型 / AI 平台", "商汤、MiniMax、阶跃星辰、稀宇科技（月之暗面上海侧）", "算法助理、数据、产品实习", "P0"],
        ["AI 芯片 / 算力", "壁仞科技、沐曦、天数智芯、芯原", "文档/测试/工具链支持", "P1"],
        ["计算机视觉 / 数据", "依图、云从（上海业务）、星环科技", "标注质检、应用交付助理", "P1"],
        ["具身智能整机", "智元机器人、傅利叶、它石智航、钛虎、开普勒", "数据采、测试、应用 demo", "P0"],
        ["服务 / 工业机器人", "擎朗智能、节卡、非夕、仙工智能、新时达", "应用工程、现场支持助理", "P0"],
        ["自动驾驶 / 出行 AI", "魔视智能、木蚁机器人、相关上海研发中心", "数据闭环、仿真支持", "P1"],
    ]
    make_table(s, Inches(0.4), Inches(1.5), Inches(12.5), Inches(5.0),
               headers, rows, font_size=11)

    # ── 4 上海 AI 企业名单 ──
    s, p = new_slide()
    slide_header(s, "一、上海人工智能企业名录（推送池 A）", "论证用短名单：落地以实际签约状态为准", p, TOTAL)
    headers = ["企业", "赛道", "适合推送方向", "合作难度", "预估通道成本(元/人)"]
    rows = [
        ["商汤科技", "视觉大模型", "算法应用 / 产品助理 / 数据", "中高", "3,200–4,000"],
        ["MiniMax", "大模型", "提示工程 / 评测 / 内容应用", "中高", "3,000–3,800"],
        ["阶跃星辰", "大模型", "应用开发助理 / 评测标注", "中", "2,800–3,500"],
        ["壁仞科技", "AI 芯片", "工具文档 / 测试支持", "高", "3,000–3,800"],
        ["沐曦集成电路", "GPU", "软件栈文档 / 适配测试", "高", "3,000–3,800"],
        ["天数智芯", "AI 芯片", "资料研究 / 工具支持", "中高", "2,800–3,500"],
        ["芯原股份", "芯片设计服务", "IP 文档 / 流程助理", "中", "2,500–3,200"],
        ["星环科技", "大数据 / AI", "数据平台实施助理", "中", "2,200–2,800"],
        ["依图科技", "视觉 AI", "方案支持 / 标注质检", "中", "2,200–2,800"],
        ["泛微网络", "AI 办公", "实施 / 客户成功助理", "低", "1,500–2,200"],
    ]
    make_table(s, Inches(0.4), Inches(1.45), Inches(12.5), Inches(5.2),
               headers, rows, font_size=11)

    # ── 5 具身智能企业名单 ──
    s, p = new_slide()
    slide_header(s, "一、具身智能 / 机器人企业名录（推送池 B）", "上海本地为主，履历辨识度高，建议作为首批签约重点", p, TOTAL)
    headers = ["企业", "赛道", "适合推送方向", "合作难度", "预估通道成本(元/人)"]
    rows = [
        ["智元机器人", "人形 / 具身智能", "数据采集、测试、应用 demo", "中高", "3,200–4,200"],
        ["傅利叶智能", "康复具身智能", "产品助理、临床支持助理", "中", "2,500–3,200"],
        ["它石智航", "工业具身智能", "现场应用、数据标注", "中高", "3,000–3,800"],
        ["擎朗智能", "商用服务机器人", "运营支持、场景落地助理", "中", "2,000–2,800"],
        ["节卡机器人", "协作机器人", "应用工程师助理", "中", "2,200–2,800"],
        ["非夕机器人", "柔性装配", "工艺/测试支持", "中", "2,200–2,800"],
        ["仙工智能", "工业 SLAM", "导航测试、文档支持", "中", "2,000–2,600"],
        ["开普勒机器人", "人形 / 物流", "整机测试、场景验证", "中高", "2,800–3,500"],
        ["钛虎机器人", "人形机器人", "结构/电控测试助理", "中高", "2,800–3,500"],
        ["灵初智能", "具身数据 / 模型", "数据采集、评测", "中高", "3,000–3,800"],
        ["首形科技", "仿人头部等", "硬件测试、展示支持", "中", "2,200–2,800"],
        ["新时达（机器人）", "工业机器人", "焊接/集成应用助理", "中", "2,000–2,600"],
    ]
    make_table(s, Inches(0.35), Inches(1.4), Inches(12.6), Inches(5.3),
               headers, rows, font_size=10)

    # ── 6 收费项目总表 ──
    s, p = new_slide()
    slide_header(s, "二、明确收费项目（学生端价目）", "收什么、收多少、含什么 —— 对外报价以此表为准", p, TOTAL)
    headers = ["收费项目", "标准价(元/人)", "包含", "不包含", "适用人群"]
    rows = [
        ["线上实习·证明版", "199", "任务+编号证明", "企业驻场、推荐信", "高中生/体验"],
        ["证明+申报支持", "699", "证明+申报邮箱+材料清单", "推荐信深度润色", "高中申请"],
        ["证明+推荐信协助", "1,680", "模板+润色+导师签字流", "包录取", "海外申请"],
        ["科技/具身智能项目实习", "5,480", "4–6周项目+周报+证明", "大厂品牌背书", "大学生"],
        ["上海 AI / 头部通道营", "11,800", "品牌项目+导师+答辩", "包转正/包录取", "大学生优生"],
        ["毕业生就业转化包", "9,800", "过渡实习+内推辅导+复盘", "包 Offer", "毕业生"],
    ]
    make_table(s, Inches(0.35), Inches(1.45), Inches(12.6), Inches(4.6),
               headers, rows, font_size=11)
    add_text(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.5),
             "支付节奏：定金 30% → 开营/推送确认 40% → 结业出具证明 30%。学校/机构团购：标准价 6–8 折。",
             size=13, color=DARK)

    # ── 7 推送成本 ──
    s, p = new_slide()
    slide_header(s, "二、推送至企业的成本（我们花多少）", "「收了学生再推过去」—— 单人通道成本与毛利对照", p, TOTAL)
    headers = ["产品档", "学生收费", "企业通道成本", "运营/导师", "获客分摊", "单人总成本", "单人毛利"]
    rows = [
        ["证明版（公益）", "199", "0–50", "60", "40", "约 150", "约 50"],
        ["申报/推荐信", "699–1,680", "0–100", "200–700", "80–150", "220–950", "约 50%"],
        ["具身/科技项目", "5,480", "2,000–2,800", "600", "900", "约 3,900–4,500", "约 1,000–1,600"],
        ["上海 AI 通道营", "11,800", "3,000–4,200", "900", "1,200", "约 5,500–6,800", "约 5,000–6,300"],
        ["就业转化包", "9,800", "2,200–3,200", "1,200", "1,000", "约 4,800–5,800", "约 4,000–5,000"],
    ]
    make_table(s, Inches(0.3), Inches(1.45), Inches(12.7), Inches(3.8),
               headers, rows, font_size=11)
    add_bullet_list(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.4),
                    [
                        "企业通道成本 = 与企业/课题方按人头或批次结算的费用（或内部导师课题等价成本）",
                        "首批上海具身智能企业：建议按 2,000–3,200 元/人谈价；AI 头部通道按 3,000–4,200 元/人预留",
                        "未签约前不得对外承诺具体企业名称；签约后按「企业匹配表」一对一推送",
                    ],
                    size=13, color=DARK)

    # ── 8 成本构成图示 ──
    s, p = new_slide()
    slide_header(s, "二、单人成本结构示意（以「具身/科技项目 5,480」为例）", "让财务与教务对齐：钱花在哪、利润留在哪", p, TOTAL)
    items = [
        ("企业通道", "2,400", "约 44%", ORANGE),
        ("导师批改", "600", "约 11%", ACCENT),
        ("获客 CAC", "900", "约 16%", PRIMARY),
        ("运营客服", "450", "约 8%", GREY),
        ("平台合规", "250", "约 5%", GREY),
        ("毛利空间", "约 1,880", "约 34%", GOLD),
    ]
    for i, (name, amt, pct, color) in enumerate(items):
        x = 0.5 + (i % 3) * 4.2
        y = 1.6 + (i // 3) * 2.5
        add_rounded(s, Inches(x), Inches(y), Inches(3.9), Inches(2.2), fill=LIGHT)
        add_rect(s, Inches(x), Inches(y), Inches(3.9), Inches(0.12), fill=color)
        add_text(s, Inches(x + 0.2), Inches(y + 0.4), Inches(3.5), Inches(0.4),
                 name, size=16, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
        add_text(s, Inches(x + 0.2), Inches(y + 0.95), Inches(3.5), Inches(0.5),
                 f"{amt} 元", size=28, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(s, Inches(x + 0.2), Inches(y + 1.55), Inches(3.5), Inches(0.35),
                 pct, size=14, color=GREY, align=PP_ALIGN.CENTER)

    # ── 9 总流程 ──
    s, p = new_slide()
    slide_header(s, "三、端到端流程（收生 → 推企业 → 出证明）", "七步闭环，每步有责任人与时间要求", p, TOTAL)

    steps = [
        ("01", "咨询报名", ["定档套餐", "签协议收定金"]),
        ("02", "背景评估", ["专业/意向匹配", "出个人推送方案"]),
        ("03", "企业匹配", ["从名录选 1–2 家", "确认名额与课题"]),
        ("04", "推送上岗", ["提交简历材料", "企业确认接收"]),
        ("05", "实习执行", ["周报打卡", "导师轻量辅导"]),
        ("06", "考核结业", ["成果验收", "企业评语"]),
        ("07", "出具交付", ["编号证明", "推荐信/申报可选"]),
    ]
    for i, (num, title, lines) in enumerate(steps):
        x = 0.35 + i * 1.85
        flow_box(s, Inches(x), Inches(1.55), Inches(1.7), Inches(2.6),
                 num, title, lines, fill=PRIMARY if i % 2 == 0 else ACCENT)
        if i < 6:
            arrow_right(s, Inches(x + 1.68), Inches(2.6))

    headers = ["步骤", "责任角色", "SLA（工作时限）", "完成标志"]
    rows = [
        ["01–02 报名评估", "招生顾问 + 教务", "报名后 2 个工作日", "匹配方案确认书"],
        ["03–04 匹配推送", "企业对接人", "方案确认后 5 个工作日", "企业接收确认函"],
        ["05 实习执行", "学员 + 企业导师 + 教务", "按 4–6 周计划", "周报齐全"],
        ["06–07 结业交付", "教务 + 合规", "结业后 3 个工作日出具证明", "证明编号可查"],
    ]
    make_table(s, Inches(0.4), Inches(4.4), Inches(12.5), Inches(2.5),
               headers, rows, font_size=11)

    # ── 10 流程细节：谁对接谁 ──
    s, p = new_slide()
    slide_header(s, "三、推送作业细则（教务 / 企业对接人手册摘要）", "避免「收了学生不知道往哪推」", p, TOTAL)
    make_card(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.2),
              "匹配规则", [
                  "高中生 → 优先公益证明 / 轻量远程课题",
                  "大学生理工 → 具身智能 / 机器人 / AI 芯片",
                  "大学生综合 → AI 应用 / SaaS / 数据标注质检",
                  "毕业生 → 就业包：可转正意向企业优先",
                  "每人主推 1 家，备选 1 家（防履约失败）",
                  "未签约企业不得写入协议承诺条款",
              ], accent=PRIMARY, body_size=13)
    make_card(s, Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.2),
              "推送材料包（一次备齐）", [
                  "学员简历（统一模板）",
                  "意向岗位与可实习时段",
                  "家长/本人知情同意（未成年）",
                  "服务协议与保密承诺",
                  "企业侧：接收确认 + 课题说明 + 导师联系人",
                  "教务侧：开营通知与周报链接",
              ], accent=ACCENT, body_size=13)

    # ── 11 资金流与角色 ──
    s, p = new_slide()
    slide_header(s, "三、资金流与角色分工", "学生付钱 → 我们结算企业 → 交付证明", p, TOTAL)
    headers = ["角色", "做什么", "关键产出", "关联费用"]
    rows = [
        ["招生顾问", "咨询、定档、收款、签约", "报名表、定金到账", "获客/提成（计入 CAC）"],
        ["教务运营", "评估、排期、周报督学、出证明", "匹配方案、证明编号", "运营分摊"],
        ["企业对接人", "维护名录、谈价、推送、履约跟进", "企业确认函、结算单", "通道成本"],
        ["企业导师", "布置任务、评语、可选推荐信", "周评/结业评语", "含在通道或导师费"],
        ["合规/财务", "协议审核、开票、分阶段结算", "回款与成本台账", "平台合规成本"],
    ]
    make_table(s, Inches(0.4), Inches(1.5), Inches(12.5), Inches(4.0),
               headers, rows, font_size=12)
    add_text(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.8),
             [
                 "结算建议：企业通道费在「企业确认接收」后预付 50%，学员结业验收后结清 50%。",
                 "若企业临时无法接收：启用备选企业；仍无法交付则按协议退还未履行部分费用。",
             ],
             size=13, color=DARK)

    # ── 12 一页纸结论 ──
    s, p = new_slide()
    slide_header(s, "一页纸结论（可直接对齐胡老师问题）", "名单有了 · 收费清了 · 成本清了 · 流程清了", p, TOTAL)
    make_card(s, Inches(0.45), Inches(1.45), Inches(4.0), Inches(5.2),
              "企业名单", [
                  "推送池 A：上海 AI 约 10 家起",
                  "推送池 B：具身智能约 12 家起",
                  "P0 先签：智元/傅利叶/擎朗/节卡/商汤等",
                  "状态字段：未开始→对接中→已签约",
                  "详见本册第 4–5 页",
              ], accent=PRIMARY, body_size=13)
    make_card(s, Inches(4.65), Inches(1.45), Inches(4.0), Inches(5.2),
              "收费与推送成本", [
                  "学生收费：199 → 11,800 分档",
                  "主推商业档标准价 5,480",
                  "推企业成本约 2,000–4,200/人",
                  "商业档单人毛利约 1,000–6,000",
                  "详见本册第 6–8 页",
              ], accent=ACCENT, body_size=13)
    make_card(s, Inches(8.85), Inches(1.45), Inches(4.0), Inches(5.2),
              "流程", [
                  "报名→评估→匹配→推送",
                  "→实习→结业→出证明",
                  "匹配 2 日 / 推送 5 日 / 证明 3 日",
                  "一人主推+一人备选",
                  "详见本册第 9–11 页",
              ], accent=ORANGE, body_size=13)

    out = Path(__file__).resolve().parents[1] / "exports"
    out.mkdir(parents=True, exist_ok=True)
    path = out / "实习推送_企业名录与收费流程说明.pptx"
    prs.save(str(path))
    print(f"已生成: {path}")
    return path


if __name__ == "__main__":
    main()
