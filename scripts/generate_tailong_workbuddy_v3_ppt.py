# -*- coding: utf-8 -*-
"""生成《泰隆银行 × 腾讯云 Work Buddy 战略合作实施方案 V3.0》PPT。

在 V2.0 落地版基础上升级为「可汇报 + 可决策 + 可执行」作战包：
- 双闭环战略叙事更聚焦
- 三项决策卡与会后两周动作前置
- 90 天里程碑 / 准出门槛可视化
- 商业漏斗、投入分担、RACI 可直接用于启动会
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from office_compat import apply_pptx_compat

# ----------------------------------------------------------------------------
# 主题：银行深蓝 + 云青 + 琥珀强调（避开紫白/奶油赤陶/报纸风）
# ----------------------------------------------------------------------------
INK = RGBColor(0x0A, 0x25, 0x40)
INK_2 = RGBColor(0x12, 0x3A, 0x5C)
PANEL = RGBColor(0x16, 0x48, 0x6B)
PANEL_LT = RGBColor(0x1E, 0x5A, 0x80)
CYAN = RGBColor(0x2B, 0xB3, 0xC0)
CYAN_LT = RGBColor(0xA8, 0xE6, 0xEC)
AMBER = RGBColor(0xE8, 0xA8, 0x38)
AMBER_LT = RGBColor(0xF5, 0xD4, 0x8A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF = RGBColor(0xE8, 0xF1, 0xF5)
MUTED = RGBColor(0x8F, 0xB0, 0xC2)
OK = RGBColor(0x3D, 0xB3, 0x8A)
WARN = RGBColor(0xE0, 0x9A, 0x3C)
RISK = RGBColor(0xD9, 0x5A, 0x5A)

FONT = "微软雅黑"
SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]
PAGE = 0


def _font(run, size, color, bold=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = FONT
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT)


def slide(bg=INK):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid()
    r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def rect(s, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE, line=None):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line:
        sp.line.color.rgb = line
        sp.line.width = Pt(1.25)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def rrect(s, x, y, w, h, color, line=None, adj=0.08):
    sp = rect(s, x, y, w, h, color, MSO_SHAPE.ROUNDED_RECTANGLE, line)
    try:
        sp.adjustments[0] = adj
    except Exception:
        pass
    return sp


def txt(s, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    for i, spec in enumerate(lines):
        text, size, color, bold = spec[0], spec[1], spec[2], spec[3]
        space_after = spec[4] if len(spec) > 4 else 4
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = text
        _font(run, size, color, bold)
    return tb


def header(s, title, subtitle=None, tag=None):
    rect(s, Inches(0), Inches(0), Inches(0.14), SH, CYAN)
    txt(s, Inches(0.45), Inches(0.24), Inches(10.8), Inches(0.5),
        [(title, 24, WHITE, True)])
    if subtitle:
        txt(s, Inches(0.47), Inches(0.74), Inches(10.8), Inches(0.32),
            [(subtitle, 12, CYAN_LT, False)])
    if tag:
        chip = rrect(s, Inches(11.2), Inches(0.28), Inches(1.75), Inches(0.4), PANEL, CYAN)
        tf = chip.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = tag
        _font(r, 11, CYAN_LT, True)
    rect(s, Inches(0.45), Inches(1.15), Inches(12.4), Pt(1.5), CYAN)


def footer(s, n):
    txt(s, Inches(0.45), Inches(7.1), Inches(10.5), Inches(0.28),
        [("泰隆银行 × 腾讯云 Work Buddy · 战略合作实施方案 V3.0 作战包", 9, MUTED, False)])
    txt(s, Inches(11.9), Inches(7.1), Inches(1.1), Inches(0.28),
        [(str(n), 9, MUTED, False)], align=PP_ALIGN.RIGHT)


def new(title=None, subtitle=None, tag=None):
    global PAGE
    PAGE += 1
    s = slide()
    if title:
        header(s, title, subtitle, tag)
    footer(s, PAGE)
    return s


def card(s, x, y, w, h, title, body_lines, accent=CYAN):
    rrect(s, x, y, w, h, PANEL, adj=0.06)
    rect(s, x, y, Inches(0.1), h, accent)
    txt(s, x + Inches(0.22), y + Inches(0.12), w - Inches(0.35), Inches(0.35),
        [(title, 14, accent, True)])
    lines = []
    for t in body_lines:
        lines.append((t, 11, OFF, False, 6))
    txt(s, x + Inches(0.22), y + Inches(0.48), w - Inches(0.35), h - Inches(0.55), lines)


def kpi(s, x, y, w, h, num, label, sub=None):
    rrect(s, x, y, w, h, PANEL, adj=0.08)
    txt(s, x + Inches(0.15), y + Inches(0.18), w - Inches(0.3), Inches(0.45),
        [(num, 22, AMBER, True)], align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.1), y + Inches(0.65), w - Inches(0.2), Inches(0.35),
        [(label, 12, WHITE, True)], align=PP_ALIGN.CENTER)
    if sub:
        txt(s, x + Inches(0.1), y + Inches(0.98), w - Inches(0.2), Inches(0.35),
            [(sub, 10, MUTED, False)], align=PP_ALIGN.CENTER)


# ============================================================================
# 封面
# ============================================================================
s = slide()
PAGE = 1
# 氛围层
rect(s, 0, 0, SW, SH, INK)
rect(s, Inches(8.2), 0, Inches(5.2), SH, INK_2)
rect(s, Inches(0), Inches(6.85), SW, Inches(0.65), PANEL)
txt(s, Inches(0.7), Inches(1.5), Inches(7.2), Inches(0.4),
    [("战略合作实施方案  V3.0 作战包", 16, CYAN, True)])
txt(s, Inches(0.7), Inches(2.05), Inches(7.5), Inches(1.4),
    [("泰隆银行 × 腾讯云", 36, WHITE, True, 8),
     ("Work Buddy", 36, AMBER, True)])
txt(s, Inches(0.7), Inches(3.85), Inches(7.2), Inches(1.0),
    [("AI 生产力 × 跨境出海 · 双闭环生态共建", 18, CYAN_LT, False, 8),
     ("把 V2.0 的方向，翻译成下周一就能开工的动作、数字与责任", 14, OFF, False)])
for i, t in enumerate(["90 天双试点", "可计量闭环", "准出门槛", "Excel 作战表"]):
    chip = rrect(s, Inches(0.7 + i * 1.85), Inches(5.35), Inches(1.7), Inches(0.42), PANEL, CYAN)
    tf = chip.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = t
    _font(r, 11, CYAN_LT, True)
txt(s, Inches(8.55), Inches(1.8), Inches(4.3), Inches(3.8),
    [("提交对象", 12, MUTED, False, 6),
     ("浙江泰隆商业银行", 16, WHITE, True, 4),
     ("腾讯云（含 Work Buddy）", 16, WHITE, True, 14),
     ("文档定位", 12, MUTED, False, 6),
     ("高层决策会 + 联合 PMO 启动包", 14, OFF, False, 14),
     ("编制日期", 12, MUTED, False, 6),
     ("2026 年 7 月", 14, AMBER, True)])
txt(s, Inches(0.7), Inches(7.0), Inches(11), Inches(0.3),
    [("内部讨论使用 · 金额比例为测算假设 · 正式条款以双方协议为准", 10, MUTED, False)])
footer(s, PAGE)

# ============================================================================
# V2→V3 升级说明
# ============================================================================
s = new("从 V2.0 落地版到 V3.0 作战包", "不改战略方向，强化可执行性、可计量性与可决策性", "升级说明")
upgrades = [
    ("战略表达", "双主线概念完整", "决策卡前置 + 一页纸结论句 + 对外口径锁定"),
    ("商业模型", "三情景测算框架", "投入分担建议 + 漏斗敏感变量 + 止损条件可勾选"),
    ("执行计划", "W0–W13 任务表", "里程碑闸门 + 交付物验收标准 + 责任到岗"),
    ("联名卡", "说明书级参数", "权益成本封顶仪表盘 + 反套利规则可落地"),
    ("合规风控", "红黄绿灯与清单", "字段白名单草案 + 审批路径时点表"),
    ("配套交付", "PDF 沟通稿", "本 PPT + Excel 作战表，直接进入联合 PMO"),
]
y0 = Inches(1.4)
for i, (dim, v2, v3) in enumerate(upgrades):
    y = y0 + Inches(i * 0.85)
    rrect(s, Inches(0.45), y, Inches(12.4), Inches(0.75), PANEL, adj=0.05)
    txt(s, Inches(0.65), y + Inches(0.18), Inches(1.8), Inches(0.4), [(dim, 14, AMBER, True)])
    txt(s, Inches(2.6), y + Inches(0.1), Inches(4.3), Inches(0.55),
        [("V2.0", 10, MUTED, False, 2), (v2, 12, OFF, False)])
    txt(s, Inches(7.2), y + Inches(0.1), Inches(5.3), Inches(0.55),
        [("V3.0", 10, CYAN, False, 2), (v3, 12, WHITE, True)])

# ============================================================================
# 一句话方案 + 基本主张
# ============================================================================
s = new("一句话方案与基本主张", "采购是入场券，联名卡是获客工具，资产是联合经营能力", "执行摘要")
rrect(s, Inches(0.45), Inches(1.4), Inches(12.4), Inches(1.55), PANEL, adj=0.05)
txt(s, Inches(0.7), Inches(1.55), Inches(11.9), Inches(1.25),
    [("方案定位", 12, CYAN, True, 6),
     ("以战略协议为总牵引，以 Work Buddy 小规模采购与内部示范为起点，以联名信用卡权益池为获客抓手，", 13, OFF, False, 4),
     ("以 20 万+ 存量客户为增长基础，以跨境金融为高价值客群锚点，形成「内部示范→权益获客→场景转化→数据复盘→生态扩容」闭环。", 13, WHITE, False)])
card(s, Inches(0.45), Inches(3.2), Inches(6.05), Inches(3.4),
     "银行真正获得什么",
     ["• 以 AI 使用行为 + 跨境资金流观察客户的新能力",
      "• 生产力场景下的卡活跃与粘性",
      "• 科技型客群「第一联系人」位置",
      "• 可复制的 90 天验证机制，而不是一次性采购"],
     AMBER)
card(s, Inches(6.75), Inches(3.2), Inches(6.1), Inches(3.4),
     "腾讯云真正获得什么",
     ["• 有信任、有支付、有场景的小微企业入口",
      "• 客户经理网络带来的线下密度",
      "• 联名权益驱动的真实激活与续费",
      "• 银行背书下的行业案例与标杆"],
     CYAN)

# ============================================================================
# 双闭环
# ============================================================================
s = new("双闭环结构：同一批客户的两个价值切面", "主线 A 全量试点 · 主线 B 样本验证 · 共用权益与看板", "战略架构")
card(s, Inches(0.45), Inches(1.4), Inches(6.05), Inches(5.2),
     "主线 A · AI 生产力金融闭环",
     ["核心问题：想用 AI，但不会用、不敢用、算不清账",
      "银行抓手：联名卡 / 分期 / 客户经理 / 权益运营",
      "云侧抓手：Work Buddy / Agent / Token / 云折扣",
      "收入形态：卡收入、分期、分润、粘性",
      "首期目标：90 天全量跑通闭环",
      "",
      "关键验证：领取 → 激活 → 线索是否连续"],
     CYAN)
card(s, Inches(6.75), Inches(1.4), Inches(6.1), Inches(5.2),
     "主线 B · 科创跨境出海闭环",
     ["核心问题：钱怎么出去、怎么回来、合不合规",
      "银行抓手：FT 账户 / 跨境贷 / 退税 / 汇率",
      "云侧抓手：出海云 / 多语言 Agent / 合规工具",
      "收入形态：利息、结售汇、结算沉淀",
      "首期目标：货架就位 + 20–30 户样本接洽",
      "",
      "关键验证：需求摸排 → 方案 → 开户落地"],
     AMBER)

# ============================================================================
# 关键数字
# ============================================================================
s = new("首期关键数字（测算口径，非承诺）", "用于建立讨论基准；正式目标值写入补充协议", "关键数字")
metrics = [
    ("80–150 席", "Work Buddy 内部席位", "5 个试点部门"),
    ("180–320 万", "首期采购总额", "分两笔：60% / 40%"),
    ("3,000–5,000", "90 天联名卡发卡", "12 个月 8k–24k"),
    ("150–260 万", "权益成本封顶", "按核销结算"),
    ("12–40 户", "90 天企业付费", "12 个月 45–120"),
    ("20–30 户", "跨境样本接洽", "不以放款为硬指标"),
]
for i, (n, l, sub) in enumerate(metrics):
    col, row = i % 3, i // 3
    kpi(s, Inches(0.45 + col * 4.2), Inches(1.5 + row * 2.5), Inches(3.95), Inches(2.2), n, l, sub)

# ============================================================================
# 三项决策
# ============================================================================
s = new("本次会议需要的三项决策", "不做决策 = 项目停留在意向阶段", "决策卡")
decisions = [
    ("01", "是否启动 90 天双试点", "AI 主线全量 + 跨境主线样本", "否则双方无法投入实质资源", RISK),
    ("02", "是否成立联合 PMO", "双方各 1 名负责人，直报分管领导", "否则跨部门进度被日常业务挤压", AMBER),
    ("03", "是否两周内完成测算与合规预审", "商务测算表 + 合规清单作为签约依据", "否则谈判缺少数据基础、反复空转", CYAN),
]
for i, (no, title, detail, risk, color) in enumerate(decisions):
    y = Inches(1.4 + i * 1.7)
    rrect(s, Inches(0.45), y, Inches(12.4), Inches(1.5), PANEL, adj=0.05)
    rrect(s, Inches(0.7), y + Inches(0.35), Inches(1.1), Inches(0.8), color, adj=0.15)
    txt(s, Inches(0.7), y + Inches(0.5), Inches(1.1), Inches(0.5),
        [(no, 20, INK if color != RISK else WHITE, True)], align=PP_ALIGN.CENTER)
    txt(s, Inches(2.1), y + Inches(0.25), Inches(10), Inches(0.4), [(title, 18, WHITE, True)])
    txt(s, Inches(2.1), y + Inches(0.7), Inches(10), Inches(0.55),
        [(detail, 13, OFF, False, 4), (f"若不决：{risk}", 12, MUTED, False)])

# ============================================================================
# 建议结论句
# ============================================================================
s = new("建议会议结论表述", "首轮会议只锁定方向与组织，不一次性写死全年 KPI", "结论句")
rrect(s, Inches(0.45), Inches(1.6), Inches(12.4), Inches(3.2), PANEL, adj=0.05)
txt(s, Inches(0.8), Inches(1.9), Inches(11.7), Inches(2.6),
    [("建议结论（可直接写入纪要）", 14, AMBER, True, 12),
     ("双方原则同意围绕企业 AI 协作与小微企业数字化服务开展战略合作，同意启动 90 天联合试点，", 16, WHITE, False, 8),
     ("成立联合项目办公室，两周内提交项目章程、商务测算与合规清单，作为签署框架协议的依据。", 16, WHITE, False)])
card(s, Inches(0.45), Inches(5.05), Inches(12.4), Inches(1.55),
     "会后两周必须交出的四份文件",
     ["① 项目章程 v0.1（含 RACI）　　② 商务测算表 v1　　③ 数据共享正负面清单　　④ MVP 权益方案骨架"],
     CYAN)

# ============================================================================
# 三轮驱动
# ============================================================================
s = new("三轮驱动模型", "先自己用起来，再用权益带进来，最后把高价值客户留住", "增长引擎")
wheels = [
    ("A 轮 · 内部示范", "客户经理 / 信用卡 / 市场 / 知识服务 / 运营",
     ["沉淀可对外案例（脱敏）", "客户经理先成为熟练用户", "场景卡 + AI 使用规范同步上线"], CYAN),
    ("B 轮 · 客户增长", "联名卡 + 腾讯云权益包",
     ["体验 → 试用 → 付费 → 续费", "每级有定义、埋点、责任人", "反馈反哺权益与话术"], AMBER),
    ("C 轮 · 科创深耕", "具身智能 / AI / 制造 / 跨境电商",
     ["参访 + 诊断 + 跨境方案", "银行升级为综合服务集成商", "首期小而精，做标杆"], OK),
]
for i, (title, who, bullets, color) in enumerate(wheels):
    x = Inches(0.45 + i * 4.2)
    rrect(s, x, Inches(1.4), Inches(4.0), Inches(5.2), PANEL, adj=0.05)
    rect(s, x, Inches(1.4), Inches(4.0), Inches(0.12), color)
    txt(s, x + Inches(0.25), Inches(1.7), Inches(3.5), Inches(0.4), [(title, 16, color, True)])
    txt(s, x + Inches(0.25), Inches(2.2), Inches(3.5), Inches(0.7), [(who, 12, MUTED, False)])
    lines = [(f"• {b}", 13, OFF, False, 10) for b in bullets]
    txt(s, x + Inches(0.25), Inches(3.1), Inches(3.5), Inches(3.0), lines)

# ============================================================================
# 闭环八步
# ============================================================================
s = new("客户增长闭环八步", "每一步必须有埋点、定义与责任人，否则无法优化", "增长闭环")
steps = [
    ("1", "客户触达"), ("2", "申请/绑卡"), ("3", "消费达标"), ("4", "领取权益"),
    ("5", "场景体验"), ("6", "顾问跟进"), ("7", "付费/续费"), ("8", "联合复盘"),
]
for i, (n, name) in enumerate(steps):
    col, row = i % 4, i // 4
    x = Inches(0.45 + col * 3.2)
    y = Inches(1.55 + row * 2.4)
    rrect(s, x, y, Inches(3.0), Inches(2.05), PANEL, adj=0.08)
    txt(s, x + Inches(0.2), y + Inches(0.35), Inches(2.6), Inches(0.55),
        [(n, 28, AMBER if row == 0 else CYAN, True)], align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.2), y + Inches(1.1), Inches(2.6), Inches(0.5),
        [(name, 16, WHITE, True)], align=PP_ALIGN.CENTER)

# ============================================================================
# 采购与场景
# ============================================================================
s = new("模块一：采购结构与五场景试点", "小步快跑：试点结果挂钩第二笔付款与扩容", "产品采购")
# 采购档
headers = ["组成", "保守", "中性", "积极"]
rows = [
    ["席位", "80", "120", "150"],
    ["场景 Agent", "3", "5", "8"],
    ["测算总额", "≈180万", "≈240万", "≈320万"],
]
rrect(s, Inches(0.45), Inches(1.4), Inches(6.1), Inches(3.0), PANEL, adj=0.05)
txt(s, Inches(0.7), Inches(1.55), Inches(5.6), Inches(0.35), [("首期采购三档", 14, CYAN, True)])
for j, h in enumerate(headers):
    txt(s, Inches(0.7 + j * 1.4), Inches(2.05), Inches(1.3), Inches(0.3), [(h, 11, MUTED, True)])
for i, row in enumerate(rows):
    for j, v in enumerate(row):
        txt(s, Inches(0.7 + j * 1.4), Inches(2.5 + i * 0.5), Inches(1.3), Inches(0.35),
            [(v, 13, AMBER if j and i == 2 else WHITE, j == 0 or i == 2)])
txt(s, Inches(0.7), Inches(4.05), Inches(5.6), Inches(0.25),
    [("支付：签约 60% + 试点验收 40%", 11, MUTED, False)])

card(s, Inches(6.8), Inches(1.4), Inches(6.05), Inches(3.0),
     "试点禁区（写入《AI 使用规范》）",
     ["× Agent 独立做授信/定价/AML/KYC",
      "× 未授权敏感信息直输模型",
      "× 模型输出直接作为对客承诺或报送"],
     RISK)

scenes = [
    ("客户经理辅助", "纪要采纳率"),
    ("信用卡运营", "出稿周期"),
    ("市场营销", "素材产出量"),
    ("内部知识服务", "一次命中率"),
    ("管理提效", "纪要覆盖率"),
]
for i, (name, m) in enumerate(scenes):
    x = Inches(0.45 + i * 2.5)
    rrect(s, x, Inches(4.7), Inches(2.35), Inches(1.85), PANEL, adj=0.08)
    txt(s, x + Inches(0.12), Inches(4.95), Inches(2.1), Inches(0.7), [(name, 13, WHITE, True)], align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.12), Inches(5.7), Inches(2.1), Inches(0.5), [(m, 11, CYAN_LT, False)], align=PP_ALIGN.CENTER)

# ============================================================================
# 联名卡
# ============================================================================
s = new("模块二：AI 生产力联名卡", "把经营性消费转化为 AI 生产力资源", "联名卡")
tiers = [
    ("标准版", "职能岗位", "体验层", "W6 小批量"),
    ("金卡版", "小微企业主", "体验+成长", "W6 小批量"),
    ("企业主版", "对公实控人", "三层全开", "W9 定向邀约"),
]
for i, (name, who, rights, when) in enumerate(tiers):
    x = Inches(0.45 + i * 4.2)
    rrect(s, x, Inches(1.4), Inches(4.0), Inches(2.4), PANEL, adj=0.06)
    txt(s, x + Inches(0.25), Inches(1.6), Inches(3.5), Inches(0.4), [(name, 18, AMBER, True)])
    txt(s, x + Inches(0.25), Inches(2.15), Inches(3.5), Inches(1.3),
        [(f"客群：{who}", 13, OFF, False, 6),
         (f"权益：{rights}", 13, OFF, False, 6),
         (f"节奏：{when}", 13, CYAN_LT, False)])

card(s, Inches(0.45), Inches(4.05), Inches(4.0), Inches(2.5),
     "体验层", ["开卡激活 / 绑卡 / 首刷", "体验席位 + 入门 Token + 公开课", "目的：降低首次门槛"], CYAN)
card(s, Inches(4.65), Inches(4.05), Inches(4.0), Inches(2.5),
     "成长层", ["月度消费达标 / 指定场景", "Agent 包 + Token 加赠 + 折扣券", "目的：提升使用深度"], AMBER)
card(s, Inches(8.85), Inches(4.05), Inches(4.0), Inches(2.5),
     "企业层", ["企业采购 / 重点任务", "顾问诊断 + 企业试用 + 活动席位", "目的：推动付费续费"], OK)

# ============================================================================
# 权益核销时序
# ============================================================================
s = new("权益核销时序与控本机制", "优先按实际核销结算；未领取不产生成本", "权益运营")
timeline = [
    ("T+0", "达标识别", "银行生成权益资格"),
    ("T+1", "资格推送", "App 权益中心提示"),
    ("T+2", "领取激活", "兑换码 / 账号绑定"),
    ("T+3", "状态回传", "腾讯云回传看板"),
    ("月末", "汇总对账", "T+5 出账 / T+10 确认"),
]
for i, (t, a, d) in enumerate(timeline):
    x = Inches(0.45 + i * 2.5)
    rrect(s, x, Inches(1.45), Inches(2.35), Inches(2.3), PANEL, adj=0.08)
    txt(s, x + Inches(0.1), Inches(1.65), Inches(2.15), Inches(0.4), [(t, 16, AMBER, True)], align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.1), Inches(2.2), Inches(2.15), Inches(0.4), [(a, 14, WHITE, True)], align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.1), Inches(2.75), Inches(2.15), Inches(0.6), [(d, 11, MUTED, False)], align=PP_ALIGN.CENTER)

controls = [
    ("总预算封顶", "150–260 万", "季度分配，超支需 steering 审批"),
    ("单客封顶", "体验+成长合计上限", "防头部客户独占预算"),
    ("领取率监控", "35%–50% 假设", "按周监控，动态调门槛"),
    ("降档触发", "70% / 85%", "自动降档或暂停发放"),
]
for i, (t, n, d) in enumerate(controls):
    x = Inches(0.45 + i * 3.2)
    rrect(s, x, Inches(4.1), Inches(3.0), Inches(2.4), PANEL, adj=0.06)
    txt(s, x + Inches(0.15), Inches(4.3), Inches(2.7), Inches(0.35), [(t, 13, CYAN, True)])
    txt(s, x + Inches(0.15), Inches(4.75), Inches(2.7), Inches(0.4), [(n, 16, AMBER, True)])
    txt(s, x + Inches(0.15), Inches(5.35), Inches(2.7), Inches(0.8), [(d, 12, OFF, False)])

# ============================================================================
# 客户运营
# ============================================================================
s = new("模块三：20 万+ 客户分层与联合运营", "客户经理是成败变量：自己先用、讲得清楚、讲了有得", "客户运营")
segs = [
    ("高意向科技", "折扣+专业服务", "定向邀约 / 顾问诊断"),
    ("成长型小微", "低成本提效", "训练营 / Agent 体验"),
    ("传统经营者", "银行背书入口", "沙龙 / 案例教育"),
    ("个人高活跃", "日常 AI 权益", "App 任务 / 月度主题"),
    ("出海意向", "跨境+出海云", "摸排 / 一揽子方案"),
]
for i, (n, v, a) in enumerate(segs):
    x = Inches(0.45 + i * 2.5)
    rrect(s, x, Inches(1.4), Inches(2.35), Inches(2.5), PANEL, adj=0.06)
    txt(s, x + Inches(0.1), Inches(1.6), Inches(2.15), Inches(0.55), [(n, 14, AMBER, True)], align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.1), Inches(2.25), Inches(2.15), Inches(0.5), [(v, 12, WHITE, False)], align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.1), Inches(2.9), Inches(2.15), Inches(0.7), [(a, 11, MUTED, False)], align=PP_ALIGN.CENTER)

sop = [
    ("1 筛选", "银行分层触达，不出完整底表"),
    ("2 教育", "每月≥2 篇行业化内容"),
    ("3 体验", "活动后 7 天内跟进"),
    ("4 商机", "书面授权后 24h 交接"),
    ("5 成交", "归属/报价/SLA 写清"),
    ("6 续费", "到期前 30 天复盘"),
]
for i, (t, d) in enumerate(sop):
    col, row = i % 3, i // 3
    x = Inches(0.45 + col * 4.2)
    y = Inches(4.15 + row * 1.25)
    rrect(s, x, y, Inches(4.0), Inches(1.1), PANEL, adj=0.06)
    txt(s, x + Inches(0.2), y + Inches(0.15), Inches(3.6), Inches(0.3), [(t, 13, CYAN, True)])
    txt(s, x + Inches(0.2), y + Inches(0.5), Inches(3.6), Inches(0.4), [(d, 12, OFF, False)])

# ============================================================================
# 跨境主线
# ============================================================================
s = new("主线二：科创跨境出海（样本验证）", "首期不做规模幻觉：账户/结算/退税优先，信贷逐笔审议", "跨境闭环")
pkgs = [
    ("基础包", "FT 账户 + 结算 + 汇率咨询 + 政策解读", "首次出海"),
    ("成长包", "基础包 + 阳光退税 + 快收 + 融资方案 + 出海云权益", "稳定跨境业务"),
    ("护航包", "成长包 + 法税物流对接 + 专属顾问 + 生态席位", "深度出海/建厂"),
]
for i, (n, c, who) in enumerate(pkgs):
    y = Inches(1.4 + i * 1.35)
    rrect(s, Inches(0.45), y, Inches(8.2), Inches(1.2), PANEL, adj=0.05)
    txt(s, Inches(0.7), y + Inches(0.2), Inches(2.2), Inches(0.35), [(n, 16, AMBER, True)])
    txt(s, Inches(3.0), y + Inches(0.2), Inches(5.4), Inches(0.8),
        [(c, 12, OFF, False, 4), (f"适用：{who}", 11, MUTED, False)])
card(s, Inches(8.9), Inches(1.4), Inches(3.95), Inches(5.1),
     "90 天节奏",
     ["W1–W4", "需求摸排 / 货架 / SOP",
      "", "W5–W8", "20–30 户接洽 / 活动 / 伙伴",
      "", "W9–W13", "方案落地 / FT 开户 / 案例",
      "", "风险提示", "境外贷后未成熟可复制",
      "→ 非信贷优先，信贷限样本"],
     CYAN)

# ============================================================================
# 商业模型
# ============================================================================
s = new("量化商业模型：六个关键变量", "最敏感变量是权益激活率；成本最不可控是领取率", "商业模型")
vars_ = [
    ("1", "权益领取率", "40%", "决定成本与转化分母"),
    ("2", "领取后激活率", "35%", "闭环最易断裂环节"),
    ("3", "激活后线索率", "12%", "企业转化输入量"),
    ("4", "线索付费转化", "18%", "收入侧规模"),
    ("5", "企业年均消费", "1.8 万", "单客价值与回本"),
    ("6", "次年续费率", "60%", "是否可持续生意"),
]
for i, (n, name, val, why) in enumerate(vars_):
    col, row = i % 3, i // 3
    x = Inches(0.45 + col * 4.2)
    y = Inches(1.4 + row * 2.55)
    rrect(s, x, y, Inches(4.0), Inches(2.35), PANEL, adj=0.06)
    txt(s, x + Inches(0.25), y + Inches(0.25), Inches(3.5), Inches(0.35), [(f"{n}  {name}", 14, CYAN, True)])
    txt(s, x + Inches(0.25), y + Inches(0.8), Inches(3.5), Inches(0.55), [(val, 28, AMBER, True)])
    txt(s, x + Inches(0.25), y + Inches(1.55), Inches(3.5), Inches(0.5), [(why, 13, OFF, False)])

# ============================================================================
# 漏斗
# ============================================================================
s = new("中性档转化漏斗（90 天 / 12 个月）", "首期真正验证第 5、6 步：领取与真实使用", "漏斗测算")
funnel = [
    ("可触达", "20,000", "60,000"),
    ("曝光", "12,000", "36,000"),
    ("申请", "6,000", "18,000"),
    ("激活卡", "4,000", "15,000"),
    ("领取权益", "1,600", "6,000"),
    ("真实激活", "560", "2,100"),
    ("高意向线索", "67", "252"),
    ("企业付费", "12", "45"),
]
txt(s, Inches(0.5), Inches(1.35), Inches(12), Inches(0.3),
    [("环节", 11, MUTED, True), ],)
# table-like cards
for i, (name, d90, d12) in enumerate(funnel):
    x = Inches(0.35 + i * 1.6)
    h = Inches(1.2 + i * 0.35)
    y = Inches(1.75)
    rrect(s, x, y, Inches(1.5), Inches(4.6), PANEL, adj=0.08)
    # fill bar visual
    bar_h = Inches(0.35 + (7 - i) * 0.35)
    rect(s, x + Inches(0.25), Inches(5.9) - bar_h, Inches(1.0), bar_h, CYAN if i < 6 else AMBER)
    txt(s, x + Inches(0.05), y + Inches(0.15), Inches(1.4), Inches(0.7),
        [(name, 11, WHITE, True)], align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.05), y + Inches(0.9), Inches(1.4), Inches(0.9),
        [("90天", 9, MUTED, False, 2), (d90, 14, AMBER, True)], align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.05), y + Inches(1.9), Inches(1.4), Inches(0.9),
        [("12月", 9, MUTED, False, 2), (d12, 14, CYAN_LT, True)], align=PP_ALIGN.CENTER)

# ============================================================================
# 投入与止损
# ============================================================================
s = new("首期投入结构与止损条件", "90 天不以盈亏平衡为目标，而是能力投资与证伪窗口", "投入止损")
inv = [
    ("Work Buddy 采购", "180–240 万", "银行承担（建议）"),
    ("权益成本", "150–210 万", "双方共担"),
    ("市场活动", "40–70 万", "双方共担"),
    ("系统对接", "30–50 万", "双方共担"),
    ("专项人力", "折算 40–60 万", "各自承担"),
]
for i, (n, a, who) in enumerate(inv):
    y = Inches(1.4 + i * 0.72)
    rrect(s, Inches(0.45), y, Inches(7.5), Inches(0.65), PANEL, adj=0.05)
    txt(s, Inches(0.7), y + Inches(0.15), Inches(3.2), Inches(0.35), [(n, 13, WHITE, True)])
    txt(s, Inches(4.0), y + Inches(0.15), Inches(2.0), Inches(0.35), [(a, 13, AMBER, True)])
    txt(s, Inches(6.1), y + Inches(0.15), Inches(1.7), Inches(0.35), [(who, 11, MUTED, False)])
card(s, Inches(8.2), Inches(1.4), Inches(4.65), Inches(5.2),
     "止损 / 重评估触发",
     ["内部月活 < 30%",
      "权益激活率 < 20%",
      "高意向线索 < 50",
      "重大合规或投诉爆发",
      "归属/数据/结算无法一致",
      "",
      "中性档回本周期",
      "约 18–30 个月",
      "高度依赖续费与二期规模"],
     RISK)

# ============================================================================
# 90天里程碑
# ============================================================================
s = new("90 天里程碑总览", "五段闸门：启动 → 准备 → 试运行 → 放量 → 评估", "执行计划")
ms = [
    ("W0–W2", "启动", "联合 PMO / 项目章程 / 合规预审", CYAN),
    ("W3–W5", "准备", "采购签约 / 场景上线 / 使用规范", CYAN),
    ("W6–W8", "试运行", "小批量发卡 / 核销链路跑通", AMBER),
    ("W9–W11", "放量", "线索交接 / 跨境样本开户", AMBER),
    ("W12–W13", "评估", "案例沉淀 / 准出决策 / 二期方案", OK),
]
for i, (w, stage, goal, color) in enumerate(ms):
    y = Inches(1.4 + i * 1.0)
    rrect(s, Inches(0.45), y, Inches(12.4), Inches(0.88), PANEL, adj=0.05)
    rrect(s, Inches(0.65), y + Inches(0.18), Inches(1.8), Inches(0.52), color, adj=0.15)
    txt(s, Inches(0.65), y + Inches(0.28), Inches(1.8), Inches(0.35),
        [(w, 14, INK, True)], align=PP_ALIGN.CENTER)
    txt(s, Inches(2.7), y + Inches(0.25), Inches(2.0), Inches(0.4), [(stage, 16, WHITE, True)])
    txt(s, Inches(5.0), y + Inches(0.25), Inches(7.5), Inches(0.4), [(goal, 14, OFF, False)])

# ============================================================================
# 三道闸门
# ============================================================================
s = new("三道准入准出门槛", "达不到就延后或缩范围，不带风险放量", "闸门")
gates = [
    ("W5 上线关", "合规预审通过", "数据清单确认", "培训完成率 ≥90%", "不通过：延后上线", CYAN),
    ("W8 链路关", "核销成功率 ≥99%", "重大故障 0 次", "投诉可控", "不通过：修复后再放量", AMBER),
    ("W13 准出关", "六项门槛全满足", "≥3 个授权案例", "单位获客成本可控", "不通过：延长 4–8 周或调整", OK),
]
for i, (title, a, b, c, d, color) in enumerate(gates):
    x = Inches(0.45 + i * 4.2)
    rrect(s, x, Inches(1.45), Inches(4.0), Inches(5.15), PANEL, adj=0.06)
    rect(s, x, Inches(1.45), Inches(4.0), Inches(0.12), color)
    txt(s, x + Inches(0.25), Inches(1.8), Inches(3.5), Inches(0.45), [(title, 18, color, True)])
    for j, line in enumerate([a, b, c]):
        txt(s, x + Inches(0.25), Inches(2.5 + j * 0.7), Inches(3.5), Inches(0.5),
            [(f"✓  {line}", 14, OFF, False)])
    txt(s, x + Inches(0.25), Inches(5.0), Inches(3.5), Inches(1.0), [(d, 13, MUTED, False)])

# ============================================================================
# KPI
# ============================================================================
s = new("90 天参考目标与复盘机制", "目标供测算，不作为战略协议刚性承诺", "指标体系")
kpis = [
    ("内部激活率", "≥80%", "周"),
    ("内部月活率", "≥50–60%", "周"),
    ("权益领取率", "35–50%", "周"),
    ("领取后激活", "≥30–40%", "周"),
    ("高意向线索转化", "≥5–10%", "月"),
    ("跨境接洽", "20–30 户", "月"),
]
for i, (n, v, f) in enumerate(kpis):
    col, row = i % 3, i // 3
    kpi(s, Inches(0.45 + col * 4.2), Inches(1.4 + row * 2.0), Inches(3.95), Inches(1.8), v, n, f"观察频率：{f}")
txt(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.0),
    [("复盘节奏：周度联合 PMO（进度/库存/风险） · 月度业务负责人（完整漏斗/成本/对账） · 季度 steering（战略/预算/准出）",
      13, OFF, False)])

# ============================================================================
# 合规
# ============================================================================
s = new("合规硬原则与 AI 红黄绿灯", "结果出行，明细不出行", "合规风控")
rrect(s, Inches(0.45), Inches(1.4), Inches(12.4), Inches(1.3), PANEL, adj=0.05)
txt(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(0.9),
    [("一条硬性原则", 13, AMBER, True, 6),
     ("银行向腾讯云传递的是「这位客户有资格领取 B 档权益」，而不是「这位客户本月消费了 6,842 元」。W2 前由合规专项小组书面确认字段清单。",
      14, WHITE, False)])
lights = [
    ("绿灯", "可直接用，事后抽查", "公开信息整理 / 会议纪要 / 素材创意", OK),
    ("黄灯", "100% 人工复核后生效", "对客文案 / 活动规则 / 制度解读", AMBER),
    ("红灯", "禁止使用", "授信定价 / AML KYC / 监管报送 / 未授权敏感信息", RISK),
]
for i, (n, d, ex, color) in enumerate(lights):
    x = Inches(0.45 + i * 4.2)
    rrect(s, x, Inches(3.0), Inches(4.0), Inches(3.5), PANEL, adj=0.06)
    rrect(s, x + Inches(1.3), Inches(3.25), Inches(1.4), Inches(0.5), color, adj=0.2)
    txt(s, x + Inches(1.3), Inches(3.32), Inches(1.4), Inches(0.4),
        [(n, 14, INK if color != RISK else WHITE, True)], align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.25), Inches(4.0), Inches(3.5), Inches(0.5), [(d, 14, WHITE, True)])
    txt(s, x + Inches(0.25), Inches(4.6), Inches(3.5), Inches(1.4), [(ex, 13, MUTED, False)])

# ============================================================================
# 组织风险
# ============================================================================
s = new("最容易被低估的风险：组织执行", "技术有解、合规有流程；客户经理不愿意讲没有技术解法", "风险提示")
card(s, Inches(0.45), Inches(1.4), Inches(6.05), Inches(5.2),
     "四件必须同步落地的事",
     ["① 自己先用：试点客户经理完成培训与首次任务，纳入考核",
      "② 讲得清楚：3 分钟话术卡 + 一页纸材料 + 5 个场景案例",
      "③ 讲了有得：激活 / 转化 / 案例贡献分别计分",
      "④ 问得到人：答疑群 + 腾讯云专家轮值，4 小时响应",
      "",
      "建议首轮会议明确：试点部门负责人对内部月活率负责"],
     AMBER)
card(s, Inches(6.75), Inches(1.4), Inches(6.1), Inches(5.2),
     "风险登记册（摘要）",
     ["高×高：组织执行（客户经理不用不推）",
      "中×高：数据隐私 / 模型输出 / 权益合规",
      "高×中：跨境贷后与汇率合规",
      "中×中：领取率失控 / 接口故障 / 归属争议",
      "低×高：舆情（银行推销 AI / 贷款买算力）",
      "",
      "Token 一律表述为「产品使用额度/服务权益」",
      "严禁与虚拟资产、投资收益产生联想"],
     RISK)

# ============================================================================
# 会后行动
# ============================================================================
s = new("会后一页纸行动清单", "从会议结束到框架协议，最长 15 个工作日", "立即执行")
actions = [
    ("1", "指定双方项目负责人并互换联系方式", "会后 2 个工作日", "双方分管领导"),
    ("2", "联合 PMO 第一次会议，发布项目章程 v0.1", "会后 5 个工作日", "双方 PMO"),
    ("3", "银行启动合规预审与数据边界梳理", "会后 5 个工作日", "银行合规"),
    ("4", "腾讯云提交产品清单、报价与技术白皮书", "会后 5 个工作日", "腾讯云"),
    ("5", "确定试点部门与人员名单", "会后 10 个工作日", "银行 PMO"),
    ("6", "完成商务测算表 v1 与合规清单", "会后 10 个工作日", "双方"),
    ("7", "向 steering 汇报，决定是否签框架协议", "会后 15 个工作日", "双方分管领导"),
]
for i, (n, a, t, o) in enumerate(actions):
    y = Inches(1.35 + i * 0.72)
    rrect(s, Inches(0.45), y, Inches(12.4), Inches(0.65), PANEL, adj=0.05)
    rrect(s, Inches(0.6), y + Inches(0.12), Inches(0.45), Inches(0.42), CYAN if i < 4 else AMBER, adj=0.2)
    txt(s, Inches(0.6), y + Inches(0.18), Inches(0.45), Inches(0.35),
        [(n, 12, INK, True)], align=PP_ALIGN.CENTER)
    txt(s, Inches(1.25), y + Inches(0.15), Inches(6.8), Inches(0.35), [(a, 13, WHITE, False)])
    txt(s, Inches(8.2), y + Inches(0.15), Inches(2.2), Inches(0.35), [(t, 12, AMBER, True)])
    txt(s, Inches(10.5), y + Inches(0.15), Inches(2.1), Inches(0.35), [(o, 11, MUTED, False)])

# ============================================================================
# 配套 Excel
# ============================================================================
s = new("配套 Excel 作战表（落地操作系统）", "PPT 用于决策对齐；Excel 用于联合 PMO 周度运转", "配套交付")
sheets = [
    ("00 总览看板", "模块目标 / 责任 / 优先级"),
    ("01 决策与确认", "三项决策 + 12 项待确认"),
    ("02 投入预算", "分档预算 + 分担建议"),
    ("03 转化漏斗", "90天/12月三情景"),
    ("04 周计划", "W0–W13 任务与交付物"),
    ("05 RACI", "责任矩阵可勾选"),
    ("06 联名卡权益", "档位 / 成本封顶 / 反套利"),
    ("07 场景卡", "五场景输入输出指标"),
    ("08 客户运营", "分层 SOP + 经理激励"),
    ("09 跨境样本", "货架 / 评分 / 节奏"),
    ("10 KPI看板", "目标 / 实际 / 红绿灯"),
    ("11 合规清单", "字段白名单 + 审批路径"),
    ("12 风险登记", "概率影响与缓解"),
    ("13 对账结算", "核销字段与月结流程"),
]
for i, (n, d) in enumerate(sheets):
    col, row = i % 4, i // 4
    x = Inches(0.45 + col * 3.2)
    y = Inches(1.4 + row * 1.4)
    rrect(s, x, y, Inches(3.05), Inches(1.25), PANEL, adj=0.06)
    txt(s, x + Inches(0.15), y + Inches(0.2), Inches(2.75), Inches(0.35), [(n, 13, AMBER, True)])
    txt(s, x + Inches(0.15), y + Inches(0.6), Inches(2.75), Inches(0.45), [(d, 11, OFF, False)])

# ============================================================================
# 结束页
# ============================================================================
s = slide()
PAGE += 1
rect(s, 0, 0, SW, SH, INK)
rect(s, 0, Inches(2.6), SW, Inches(2.4), INK_2)
txt(s, Inches(0.8), Inches(2.9), Inches(11.5), Inches(0.5),
    [("下一步：启动 90 天双试点", 28, WHITE, True)], align=PP_ALIGN.CENTER)
txt(s, Inches(0.8), Inches(3.6), Inches(11.5), Inches(0.8),
    [("成立联合 PMO · 两周内交出章程 / 测算 / 合规清单 · 用准出门槛管理扩容",
      16, CYAN_LT, False)], align=PP_ALIGN.CENTER)
txt(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.8),
    [("泰隆银行 × 腾讯云 Work Buddy", 14, AMBER, True, 6),
     ("本文件为商务沟通与内部执行草案，所有数据为测算假设，不构成承诺；条款以正式协议为准。",
      12, MUTED, False)], align=PP_ALIGN.CENTER)
footer(s, PAGE)

# 保存
out_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "deliverables")
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, "泰隆银行x腾讯云WorkBuddy_战略合作实施方案_V3.0.pptx")
apply_pptx_compat(prs)
prs.save(out_path)
print(f"已生成: {out_path}  共 {PAGE} 页")
