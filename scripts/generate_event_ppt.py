# -*- coding: utf-8 -*-
"""生成《首场·出海东南亚 总领事商务论坛 —— 策划·执行·赞助落地》PPT。

用法：
    python3 scripts/generate_event_ppt.py [输出路径.pptx]
"""
import sys
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

import event_content as C

NAVY = RGBColor(0x0B, 0x2A, 0x4A)
NAVY2 = RGBColor(0x12, 0x3A, 0x63)
GOLD = RGBColor(0xC9, 0xA2, 0x4B)
LIGHT = RGBColor(0xF4, 0xF6, 0xF9)
GREY = RGBColor(0x5B, 0x6B, 0x7B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x29, 0x33)
FONT = "Microsoft YaHei"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_font(run, size, color=DARK, bold=False, font=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', font)


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size, color=DARK, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_font(r, size, color, bold, font)
    return tb


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def header(slide, title, idx=None, total=None):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), NAVY)
    add_rect(slide, 0, Inches(1.1), SLIDE_W, Pt(3), GOLD)
    add_text(slide, Inches(0.6), Inches(0.12), Inches(11.0), Inches(0.85),
             title, 25, WHITE, True, anchor=MSO_ANCHOR.MIDDLE)
    if idx is not None:
        add_text(slide, Inches(11.6), Inches(0.12), Inches(1.3), Inches(0.85),
                 f"{idx:02d}/{total:02d}", 12, GOLD, True,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def footer(slide):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(9), Inches(0.4),
             f"{C.PROJECT_NAME} · {C.PROJECT_SUBTITLE}", 9, GREY)


def bullets(slide, x, y, w, h, items, size=16, gap=10, color=DARK,
            marker="●", marker_color=GOLD):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.line_spacing = 1.12
        rm = p.add_run()
        rm.text = marker + "  "
        set_font(rm, size - 2, marker_color, True)
        r = p.add_run()
        r.text = it
        set_font(r, size, color)
    return tb


def add_table_slide(slide, tdata, idx, total):
    header(slide, tdata["title"], idx, total)
    headers = tdata["headers"]
    rows = tdata["rows"]
    ncol = len(headers)
    nrow = len(rows) + 1
    left = Inches(0.5)
    top = Inches(1.4)
    width = Inches(12.33)
    # 自适应行高与字号，避免多行表溢出
    avail = 5.05  # inches for the table body
    row_h = min(0.42, avail / nrow)
    cell_fs = 10.5 if nrow <= 10 else (9.5 if nrow <= 13 else 8.8)
    head_fs = 12 if nrow <= 10 else 11
    height = Inches(row_h) * nrow
    gtbl = slide.shapes.add_table(nrow, ncol, left, top, width, height).table
    # 首列略宽
    first_w = int(width * 0.16) if ncol >= 4 else int(width / ncol)
    rest = int((width - first_w) / (ncol - 1)) if ncol > 1 else int(width)
    for ci in range(ncol):
        gtbl.columns[ci].width = Emu(first_w if ci == 0 else rest)
    for ci, htext in enumerate(headers):
        cell = gtbl.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.05); cell.margin_right = Inches(0.05)
        cell.margin_top = Inches(0.01); cell.margin_bottom = Inches(0.01)
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = htext
        set_font(r, head_fs, WHITE, True)
    for ri, row in enumerate(rows, start=1):
        emph = str(row[0]) in ("合计",)
        for ci, val in enumerate(row):
            cell = gtbl.cell(ri, ci)
            cell.fill.solid()
            if emph:
                cell.fill.fore_color.rgb = GOLD
            else:
                cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.05); cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.01); cell.margin_bottom = Inches(0.01)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci else PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(val)
            set_font(r, cell_fs, DARK, emph or ci == 0)
    note = tdata.get("note")
    if note:
        ny = top + height + Inches(0.1)
        add_text(slide, Inches(0.5), ny, Inches(12.33), Inches(0.9),
                 "说明：" + note, 10, GREY)
    footer(slide)


# ---- 专用页 ----

def slide_cover(prs):
    s = blank(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, 0, Inches(4.85), SLIDE_W, Pt(3), GOLD)
    add_text(s, Inches(0.9), Inches(1.15), Inches(11.5), Inches(0.5),
             C.PROJECT_TAG, 15, GOLD, True)
    add_text(s, Inches(0.9), Inches(2.35), Inches(11.7), Inches(1.5),
             C.PROJECT_NAME, 44, WHITE, True)
    add_text(s, Inches(0.92), Inches(3.75), Inches(11.5), Inches(0.8),
             C.PROJECT_SUBTITLE, 24, LIGHT)
    add_text(s, Inches(0.92), Inches(5.1), Inches(11.7), Inches(0.6),
             "东南亚六国+日韩 · 300人 · 论坛(免费)+轻社交酒会 · " + C.EVENT_DATE, 15, GOLD, True)
    add_text(s, Inches(0.92), Inches(6.6), Inches(11.5), Inches(0.5),
             "主体：总领文化传播有限公司（筹）　|　" + C.VERSION, 12, GREY)
    return s


def slide_date(prs, idx, total):
    s = blank(prs)
    header(s, "择期：黄道吉日 2026/9/28", idx, total)
    add_rect(s, Inches(0.5), Inches(1.5), Inches(5.9), Inches(2.4), NAVY)
    add_text(s, Inches(0.7), Inches(1.7), Inches(5.5), Inches(0.5), "活动日期", 16, GOLD, True)
    add_text(s, Inches(0.7), Inches(2.25), Inches(5.5), Inches(1.4),
             C.HUANGLI["date"] + "\n（周一 · 下午 13:00–18:30）", 20, WHITE, True)
    add_rect(s, Inches(6.6), Inches(1.5), Inches(6.23), Inches(2.4), LIGHT)
    add_text(s, Inches(6.8), Inches(1.65), Inches(5.9), Inches(2.2),
             "【宜】" + C.HUANGLI["yi"] + "\n\n【忌】" + C.HUANGLI["ji"] +
             "\n\n【冲】" + C.HUANGLI["chong"], 12.5, DARK)
    add_rect(s, Inches(0.5), Inches(4.15), Inches(12.33), Inches(0.6), GOLD)
    add_text(s, Inches(0.7), Inches(4.15), Inches(12), Inches(0.6),
             "择此日之由：" + C.HUANGLI["why"], 13, NAVY, True, anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, Inches(0.7), Inches(4.95), Inches(12), Inches(1.4),
            C.DATE_BACKUP + [C.HOLIDAY_NOTE], size=12.5, gap=8, marker="·")
    footer(s)
    return s


def slide_countries(prs, idx, total):
    s = blank(prs)
    header(s, "具体国家：东南亚六国 + 日韩加持", idx, total)
    # 东南亚六国 + 日韩，用色块国旗式卡片
    sea = ["新加坡", "马来西亚", "泰国", "印度尼西亚", "越南", "菲律宾"]
    ea = ["日本", "韩国"]
    add_text(s, Inches(0.6), Inches(1.35), Inches(12), Inches(0.4),
             "首场主体 · 东南亚六国", 15, NAVY, True)
    bw = Inches(1.9); gap = Inches(0.13); startx = Inches(0.6); top = Inches(1.85)
    for i, c in enumerate(sea):
        x = startx + i * (bw + gap)
        add_rect(s, x, top, bw, Inches(0.85), NAVY)
        add_text(s, x, top, bw, Inches(0.85), c, 16, WHITE, True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.6), Inches(3.0), Inches(12), Inches(0.4),
             "加持提升规格 · 东亚", 15, NAVY, True)
    for i, c in enumerate(ea):
        x = startx + i * (bw + gap)
        add_rect(s, x, Inches(3.5), bw, Inches(0.85), GOLD)
        add_text(s, x, Inches(3.5), bw, Inches(0.85), c, 16, NAVY, True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(0.5), Inches(4.7), Inches(12.33), Inches(1.7), LIGHT)
    add_text(s, Inches(0.7), Inches(4.82), Inches(12), Inches(0.4), "国别推介重点", 14, NAVY, True)
    add_text(s, Inches(0.7), Inches(5.25), Inches(12.0), Inches(1.1),
             "新加坡（总部/金融/转口）· 马来西亚（电子/清真/数据中心）· 泰国（汽车与新能源配套）"
             "\n印尼（资源/消费/数字经济）· 越南（电子/纺织/制造转移）· 菲律宾（消费/BPO/基建）"
             "\n日本（高端制造/消费/科技）· 韩国（半导体/美妆/文化科技）", 12.5, DARK)
    footer(s)
    return s


def slide_section(prs, no, title, subtitle=""):
    s = blank(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, Inches(0.9), Inches(2.6), Inches(1.6), Pt(4), GOLD)
    add_text(s, Inches(0.9), Inches(2.9), Inches(3), Inches(2), no, 96, GOLD, True)
    add_text(s, Inches(3.4), Inches(3.05), Inches(9), Inches(1.2), title, 38, WHITE,
             True, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(s, Inches(3.45), Inches(4.25), Inches(9), Inches(1), subtitle, 17, LIGHT)
    return s


def slide_goal(prs, idx, total):
    s = blank(prs)
    header(s, "一、目标与定位", idx, total)
    add_text(s, Inches(0.6), Inches(1.4), Inches(6), Inches(0.5), "三大目标", 18, NAVY, True)
    bullets(s, Inches(0.6), Inches(1.95), Inches(6.0), Inches(2.5), C.OBJECTIVES,
            size=15, gap=12)
    add_rect(s, Inches(6.95), Inches(1.45), Inches(5.85), Inches(5.0), LIGHT)
    add_text(s, Inches(7.2), Inches(1.6), Inches(5.4), Inches(0.5), "首场定位", 18, NAVY, True)
    bullets(s, Inches(7.2), Inches(2.2), Inches(5.4), Inches(4.2), C.POSITIONING,
            size=14, gap=12)
    footer(s)
    return s


def slide_model(prs, idx, total):
    s = blank(prs)
    header(s, "盈利逻辑：大活动拉赞助、小活动收费", idx, total)
    bullets(s, Inches(0.7), Inches(1.6), Inches(12), Inches(4.8), C.KEY_MODEL,
            size=17, gap=18)
    add_rect(s, Inches(0.5), Inches(6.2), Inches(12.33), Inches(0.7), GOLD)
    add_text(s, Inches(0.7), Inches(6.2), Inches(12), Inches(0.7),
             "一句话：下午论坛免费做大基数 → B 端赞助 + 展位增值 + 会客厅挂牌 → 覆盖成本并盈利。",
             13.5, NAVY, True, anchor=MSO_ANCHOR.MIDDLE)
    footer(s)
    return s


def slide_brand(prs, idx, total):
    s = blank(prs)
    header(s, "复旦品牌价值保护", idx, total)
    bullets(s, Inches(0.7), Inches(1.55), Inches(12), Inches(5.0), C.BRAND_PROTECT,
            size=15, gap=14)
    footer(s)
    return s


def slide_platform(prs, idx, total):
    s = blank(prs)
    header(s, "双平台价值最大化（复旦 × CGC）", idx, total)
    cardw = Inches(4.0)
    gap = Inches(0.27)
    startx = Inches(0.5)
    top = Inches(1.55)
    h = Inches(4.9)
    for i, (name, desc) in enumerate(C.PLATFORM_VALUE):
        x = startx + i * (cardw + gap)
        add_rect(s, x, top, cardw, h, WHITE, line=GOLD)
        add_rect(s, x, top, cardw, Inches(0.9), NAVY if i < 2 else GOLD)
        add_text(s, x + Inches(0.15), top + Inches(0.05), cardw - Inches(0.3),
                 Inches(0.8), name, 16, WHITE if i < 2 else NAVY, True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.25), top + Inches(1.15), cardw - Inches(0.5),
                 Inches(3.5), desc, 13.5, DARK)
    footer(s)
    return s


def slide_feasibility(prs, idx, total):
    s = blank(prs)
    header(s, "可行性评估（务实）", idx, total)
    bullets(s, Inches(0.7), Inches(1.55), Inches(12), Inches(5.0), C.FEASIBILITY,
            size=15.5, gap=14)
    footer(s)
    return s


def slide_next(prs, idx, total):
    s = blank(prs)
    header(s, "下一步行动", idx, total)
    top = Inches(1.7)
    rowh = Inches(1.0)
    for i, step in enumerate(C.NEXT_STEPS):
        y = top + i * (rowh + Inches(0.12))
        add_rect(s, Inches(0.6), y, Inches(0.85), rowh, GOLD)
        add_text(s, Inches(0.6), y, Inches(0.85), rowh, str(i + 1), 24, NAVY, True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(1.6), y, Inches(11.1), rowh, LIGHT)
        add_text(s, Inches(1.85), y, Inches(10.7), rowh, step, 15, DARK,
                 anchor=MSO_ANCHOR.MIDDLE)
    footer(s)
    return s


def slide_end(prs):
    s = blank(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, Inches(0.9), Inches(3.5), Inches(2.0), Pt(4), GOLD)
    add_text(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.0),
             "先跑通模式，再放大规模", 38, WHITE, True)
    add_text(s, Inches(0.92), Inches(3.8), Inches(11.5), Inches(0.8),
             "以领事为高位、以复旦为背书，把首场做成可复制的正向现金流样板。",
             17, LIGHT)
    add_text(s, Inches(0.92), Inches(6.4), Inches(11.5), Inches(0.6),
             C.PROJECT_NAME + " · " + C.PROJECT_TAG + " · " + C.VERSION, 12, GOLD)
    return s


def build(path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    total = 18
    n = [0]

    def nxt():
        n[0] += 1
        return n[0]

    slide_cover(prs)
    slide_section(prs, "01", "目标 · 定位 · 择期", "跑通正向现金流 · 沉淀核心理事 · 黄道吉日 9/28")
    slide_goal(prs, nxt(), total)
    add_table_slide(blank(prs), C.TBL_OVERVIEW, nxt(), total)
    slide_date(prs, nxt(), total)
    slide_countries(prs, nxt(), total)
    slide_section(prs, "02", "活动方案与议程", "论坛(免费) + 轻社交酒会（不做晚宴） · 详细议程")
    add_table_slide(blank(prs), C.TBL_RUNDOWN, nxt(), total)
    slide_model(prs, nxt(), total)
    slide_section(prs, "03", "收费与赞助落地", "下午免费 · B 端赞助 · 展位增值 · 会客厅 · 盈亏")
    add_table_slide(blank(prs), C.TBL_ENROLL, nxt(), total)
    add_table_slide(blank(prs), C.TBL_B_SPONSOR, nxt(), total)
    add_table_slide(blank(prs), C.TBL_ADDON, nxt(), total)
    add_table_slide(blank(prs), C.TBL_SALON_GOV, nxt(), total)
    add_table_slide(blank(prs), C.TBL_COST, nxt(), total)
    add_table_slide(blank(prs), C.TBL_PNL, nxt(), total)
    slide_section(prs, "04", "保障 · 分工 · 落地", "复旦品牌保护 · 双平台价值 · 分工 · 倒排 · 可行性")
    slide_brand(prs, nxt(), total)
    slide_platform(prs, nxt(), total)
    add_table_slide(blank(prs), C.TBL_RACI, nxt(), total)
    add_table_slide(blank(prs), C.TBL_TIMELINE, nxt(), total)
    slide_feasibility(prs, nxt(), total)
    slide_next(prs, nxt(), total)
    slide_end(prs)

    os.makedirs(os.path.dirname(path), exist_ok=True)
    prs.save(path)
    print(f"已生成 PPT：{path}（共 {len(prs.slides)} 页）")


if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else \
        "output/首场出海论坛-策划执行与赞助落地.pptx"
    build(out)
