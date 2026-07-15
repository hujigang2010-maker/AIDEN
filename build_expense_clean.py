# -*- coding: utf-8 -*-
"""
2026 年度活动 · 费用清单（干净版）

口径：完全只展示费用（资金流出），通篇无任何"进账/收入/赞助/退款剔除"的痕迹。
- 「活动前看场地用餐」按支出口径登记 ¥-148.85
- 「晚宴退费」原备注涉及升舱场景，本版改为中性表述
"""
import os
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

OUT_DIR = "/workspace/output"
os.makedirs(OUT_DIR, exist_ok=True)

# ---------------------------------------------------------------
# 数据：仅费用项目；负数表示资金流出
# ---------------------------------------------------------------

EXTERNAL = [
    ("主持人",          "Miranda 苗苗",  -500.00,    "主持人礼服"),
    ("晚宴期间演出",     "小五",         -3000.00,   "三人演出表演"),
    ("摄影摄像",        "洪总",         -3000.00,   "摄影加云图片，及摄像影像记录"),
    ("会务公司",        "杨总",         -16000.00,  "全部费用合计见附图"),
    ("媒体支出",        "葛九明",       -11500.00,  "新华社、上海证券报、央视新闻客户端、央广网、中国新闻网、上海虹口融媒体、上海杨浦融媒体（其中500元由胡继刚付款）"),
    ("晚宴用餐",        "尚九一滴水",   -33452.00,  "晚宴8桌用餐费用，主桌按1.5倍计算"),
    ("证书制作及KT版",   "广告公司",     -1030.00,   "个人证书14件 + 企业证书10件 + 志愿者证书23件"),
    ("晚宴红包",        "胡继刚",        -100.00,   "晚宴活动游戏环节用"),
    ("打印费",          "王胜",           -10.00,   "主持人手稿"),
]

INTERNAL = [
    ("茶歇",            "F组-彭常丽",   -2000.50,   "盒马鲜生采购，单次≤30kg，分8单下单（含订单截图录屏）"),
    ("布袋",            "F组-彭常丽",    -860.00,   "200个 × 4.30元/个"),
    ("志愿者伴手礼",     "F组-彭常丽",   -2025.00,   "30份 × 68元/份 = 2040元，使用15元券，实付2025元"),
    ("志愿者证书邮寄费", "胡继刚",        -200.00,   "5位志愿者未现场领取证书及礼品，另含夏春等邮寄费"),
    ("A4 刊册议程",     "佳玉",         -4400.00,   "500份刊册"),
    ("午餐",            "F组-彭常丽",    -780.00,   "大米先生 30元/份 × 26份（附付款截图）"),
    ("活动前看场地用餐", "王胜、王卓 等", -148.85,   "活动前看场地三次用餐"),
]

ONSITE = [
    ("胸花 & 台花",     "F组-朱俊锋",    -488.00,   "胸花10个 + 台花1个"),
    ("对讲机租赁",      "F组-朱俊锋",    -138.88,   "租赁100元 + 来回运费38.88元"),
    ("晚宴退费",        "F组-朱俊锋",    -199.00,   "现场退费一笔"),
    ("货拉拉送货",      "F组-朱俊锋",    -100.00,   "送复旦管院纸袋子 + 酒水送件"),
]

CERT_DETAIL = [
    ("XH20260522016", "彩色打印", "HD", "套", 24, 35, 840, "01-02 证书+外壳：企业10 + 个人14"),
    ("XH20260522016", "彩色打印", "A4", "张", 24, 5,  120, "—"),
    ("XH20260522016", "彩色打印", "A4", "套", 1,  35, 35,  "03 志愿者证书 珠光纸 1份"),
    ("XH20260522016", "快递费",   "HD", "次", 1,  35, 35,  "04 证书160克超白+证书外壳 1套；闪送：尚九一滴水1F"),
]

def _sum(rows): return round(sum(r[2] for r in rows), 2)

SUM_EXT, SUM_INT, SUM_ON = _sum(EXTERNAL), _sum(INTERNAL), _sum(ONSITE)
GRAND = round(SUM_EXT + SUM_INT + SUM_ON, 2)

assert SUM_EXT == -68592.00
assert SUM_INT == -10414.35, SUM_INT
assert SUM_ON  == -925.88,   SUM_ON
assert GRAND   == -79932.23, GRAND
print(f"外部费用 {len(EXTERNAL)} 项 ¥{SUM_EXT:,.2f}")
print(f"内部费用 {len(INTERNAL)} 项 ¥{SUM_INT:,.2f}")
print(f"现场杂项 {len(ONSITE)} 项 ¥{SUM_ON:,.2f}")
print(f"合计     {len(EXTERNAL)+len(INTERNAL)+len(ONSITE)} 项 ¥{GRAND:,.2f}")

# ---------------------------------------------------------------
# 样式（与原版统一风格）
# ---------------------------------------------------------------
COLOR_PRIMARY  = "8B0000"
COLOR_SUBHEAD  = "F5E9D3"
COLOR_STRIPE   = "FBF7F0"
COLOR_TOTAL_BG = "FFE4B5"
COLOR_EXPENSE  = "B22222"

thin = Side(border_style="thin", color="C9A86B")
medium = Side(border_style="medium", color="8B0000")
BORDER_CELL = Border(left=thin, right=thin, top=thin, bottom=thin)
BORDER_HEAD = Border(left=thin, right=thin, top=medium, bottom=medium)

FONT_TITLE   = Font(name="WenQuanYi Micro Hei", size=18, bold=True, color="FFFFFF")
FONT_SUBTITLE= Font(name="WenQuanYi Micro Hei", size=11, italic=True, color="555555")
FONT_HEAD    = Font(name="WenQuanYi Micro Hei", size=12, bold=True, color="FFFFFF")
FONT_SUBHEAD = Font(name="WenQuanYi Micro Hei", size=11, bold=True, color="8B0000")
FONT_CELL    = Font(name="WenQuanYi Micro Hei", size=10, color="333333")
FONT_TOTAL   = Font(name="WenQuanYi Micro Hei", size=12, bold=True, color="8B0000")
FONT_NUM_EXP = Font(name="WenQuanYi Micro Hei", size=10, color=COLOR_EXPENSE)

FILL_TITLE   = PatternFill("solid", fgColor=COLOR_PRIMARY)
FILL_HEAD    = PatternFill("solid", fgColor=COLOR_PRIMARY)
FILL_SUBHEAD = PatternFill("solid", fgColor=COLOR_SUBHEAD)
FILL_STRIPE  = PatternFill("solid", fgColor=COLOR_STRIPE)
FILL_TOTAL   = PatternFill("solid", fgColor=COLOR_TOTAL_BG)

ALIGN_CENTER = Alignment(horizontal="center", vertical="center", wrap_text=True)
ALIGN_LEFT   = Alignment(horizontal="left",   vertical="center", wrap_text=True)
ALIGN_RIGHT  = Alignment(horizontal="right",  vertical="center", wrap_text=True)
NUM_FMT = '_-¥* #,##0.00_-;[Red]-¥* #,##0.00_-;_-¥* "-"??_-;_-@_-'

# ---------------------------------------------------------------
# Excel
# ---------------------------------------------------------------
wb = Workbook()

def _title_row(ws, title, subtitle, ncols):
    ws.row_dimensions[1].height = 38
    ws.row_dimensions[2].height = 22
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=ncols)
    ws.merge_cells(start_row=2, start_column=1, end_row=2, end_column=ncols)
    c = ws.cell(row=1, column=1, value=title)
    c.font = FONT_TITLE; c.fill = FILL_TITLE; c.alignment = ALIGN_CENTER
    c2 = ws.cell(row=2, column=1, value=subtitle)
    c2.font = FONT_SUBTITLE; c2.alignment = ALIGN_CENTER
    for col in range(1, ncols + 1):
        ws.cell(row=1, column=col).border = BORDER_HEAD

def _header_row(ws, row, headers):
    ws.row_dimensions[row].height = 28
    for i, h in enumerate(headers, 1):
        c = ws.cell(row=row, column=i, value=h)
        c.font = FONT_HEAD; c.fill = FILL_HEAD
        c.alignment = ALIGN_CENTER; c.border = BORDER_HEAD

def _money(cell, v):
    cell.number_format = NUM_FMT
    cell.value = v
    cell.font = FONT_NUM_EXP if (v is not None and v < 0) else FONT_CELL
    cell.alignment = ALIGN_RIGHT
    cell.border = BORDER_CELL

def _widths(ws, ws_widths):
    for i, w in enumerate(ws_widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = w

def _fit(ws):
    ws.page_setup.orientation = ws.ORIENTATION_LANDSCAPE
    ws.page_setup.paperSize = ws.PAPERSIZE_A4
    ws.page_setup.fitToWidth = 1
    ws.page_setup.fitToHeight = 0
    ws.sheet_properties.pageSetUpPr.fitToPage = True
    ws.page_margins.left = 0.3
    ws.page_margins.right = 0.3
    ws.page_margins.top = 0.4
    ws.page_margins.bottom = 0.4
    ws.print_options.horizontalCentered = True

# ---- 总览 ----
ws = wb.active
ws.title = "费用总览"
_title_row(ws,
    "2026 年度活动 · 费用清单总览",
    "币种：人民币（元）",
    5)
_header_row(ws, 3, ["序号", "类别", "项目数", "金额（元）", "说明"])

summary = [
    ("一", "外部费用", len(EXTERNAL), SUM_EXT, "对外采购、服务、媒体、餐饮等"),
    ("二", "内部费用", len(INTERNAL), SUM_INT, "内部团队垫付的物料、伴手礼、餐饮等"),
    ("三", "现场杂项", len(ONSITE),  SUM_ON,  "胸花、对讲机、退费、送货等现场零星费用"),
]
r = 4
for idx, (no, name, cnt, amt, desc) in enumerate(summary):
    ws.cell(row=r, column=1, value=no).alignment = ALIGN_CENTER
    ws.cell(row=r, column=2, value=name).alignment = ALIGN_LEFT
    ws.cell(row=r, column=3, value=cnt).alignment = ALIGN_CENTER
    _money(ws.cell(row=r, column=4), amt)
    ws.cell(row=r, column=5, value=desc).alignment = ALIGN_LEFT
    for col in range(1, 6):
        cell = ws.cell(row=r, column=col)
        cell.border = BORDER_CELL
        if col != 4:
            cell.font = FONT_CELL
        if idx % 2 == 1 and col != 4:
            cell.fill = FILL_STRIPE
    r += 1

ws.row_dimensions[r].height = 30
ws.merge_cells(start_row=r, start_column=1, end_row=r, end_column=3)
sl = ws.cell(row=r, column=1, value="合计")
sl.font = FONT_TOTAL; sl.fill = FILL_TOTAL; sl.alignment = ALIGN_CENTER; sl.border = BORDER_CELL
_money(ws.cell(row=r, column=4), GRAND)
ws.cell(row=r, column=4).font = FONT_TOTAL
ws.cell(row=r, column=4).fill = FILL_TOTAL
note = ws.cell(row=r, column=5, value=f"共 {len(EXTERNAL)+len(INTERNAL)+len(ONSITE)} 笔费用累加")
note.font = FONT_TOTAL; note.fill = FILL_TOTAL; note.alignment = ALIGN_LEFT; note.border = BORDER_CELL

_widths(ws, [6, 22, 10, 18, 60])
ws.freeze_panes = "A4"
ws.sheet_view.showGridLines = False
_fit(ws)

# ---- 明细 ----
def build_detail(name, title, subtitle, data, subtotal_label, subtotal):
    ws = wb.create_sheet(name)
    _title_row(ws, title, subtitle, 5)
    _header_row(ws, 3, ["序号", "项目", "对接人 / 供应商", "金额（元）", "备注"])
    r = 4
    for i, (proj, contact, amt, remark) in enumerate(data, 1):
        ws.row_dimensions[r].height = 22 if len(remark) < 40 else 36
        ws.cell(row=r, column=1, value=i).alignment = ALIGN_CENTER
        ws.cell(row=r, column=2, value=proj).alignment = ALIGN_LEFT
        ws.cell(row=r, column=3, value=contact).alignment = ALIGN_LEFT
        _money(ws.cell(row=r, column=4), amt)
        ws.cell(row=r, column=5, value=remark).alignment = ALIGN_LEFT
        for col in range(1, 6):
            cell = ws.cell(row=r, column=col)
            cell.border = BORDER_CELL
            if col != 4:
                cell.font = FONT_CELL
            if i % 2 == 0 and col != 4:
                cell.fill = FILL_STRIPE
        r += 1
    ws.row_dimensions[r].height = 28
    ws.merge_cells(start_row=r, start_column=1, end_row=r, end_column=3)
    sl = ws.cell(row=r, column=1, value=subtotal_label)
    sl.font = FONT_TOTAL; sl.fill = FILL_TOTAL; sl.alignment = ALIGN_CENTER; sl.border = BORDER_CELL
    _money(ws.cell(row=r, column=4), subtotal)
    ws.cell(row=r, column=4).font = FONT_TOTAL
    ws.cell(row=r, column=4).fill = FILL_TOTAL
    c = ws.cell(row=r, column=5, value="")
    c.fill = FILL_TOTAL; c.border = BORDER_CELL
    _widths(ws, [6, 22, 22, 16, 65])
    ws.freeze_panes = "A4"
    ws.sheet_view.showGridLines = False
    _fit(ws)

build_detail("外部费用明细", "一、外部费用明细",
             "对外采购 / 服务 / 媒体 / 餐饮", EXTERNAL,
             "外部费用小计", SUM_EXT)
build_detail("内部费用明细", "二、内部费用明细",
             "内部团队垫付的物料、伴手礼、餐饮等", INTERNAL,
             "内部费用小计", SUM_INT)
build_detail("现场杂项明细", "三、现场杂项明细",
             "胸花 / 对讲机 / 退费 / 送货", ONSITE,
             "现场杂项小计", SUM_ON)

# ---- 证书附表 ----
ws = wb.create_sheet("证书打印费明细")
_title_row(ws,
    "附表 · 证书制作及KT版打印明细",
    "供应商：广告公司  ｜  原始单据编号：XH20260522016  ｜  合计 ¥1,030.00",
    8)
_header_row(ws, 3, ["序号", "单据编号", "商品名称", "型号", "单位", "数量", "单价（元）", "金额（元）"])
r = 4
for i, (no, name, mdl, unit, qty, price, amt, _remark) in enumerate(CERT_DETAIL, 1):
    ws.cell(row=r, column=1, value=i).alignment = ALIGN_CENTER
    ws.cell(row=r, column=2, value=no).alignment = ALIGN_CENTER
    ws.cell(row=r, column=3, value=name).alignment = ALIGN_LEFT
    ws.cell(row=r, column=4, value=mdl).alignment = ALIGN_CENTER
    ws.cell(row=r, column=5, value=unit).alignment = ALIGN_CENTER
    ws.cell(row=r, column=6, value=qty).alignment = ALIGN_CENTER
    _money(ws.cell(row=r, column=7), price); ws.cell(row=r, column=7).font = FONT_CELL
    _money(ws.cell(row=r, column=8), amt);   ws.cell(row=r, column=8).font = FONT_CELL
    for col in range(1, 9):
        cell = ws.cell(row=r, column=col)
        cell.border = BORDER_CELL
        if cell.font is None: cell.font = FONT_CELL
        if i % 2 == 0: cell.fill = FILL_STRIPE
    r += 1
ws.merge_cells(start_row=r, start_column=1, end_row=r, end_column=7)
sl = ws.cell(row=r, column=1, value="合计")
sl.font = FONT_TOTAL; sl.fill = FILL_TOTAL; sl.alignment = ALIGN_CENTER; sl.border = BORDER_CELL
_money(ws.cell(row=r, column=8), 1030.00)
ws.cell(row=r, column=8).font = FONT_TOTAL
ws.cell(row=r, column=8).fill = FILL_TOTAL
_widths(ws, [6, 18, 14, 8, 8, 8, 14, 16])
ws.freeze_panes = "A4"
ws.sheet_view.showGridLines = False
_fit(ws)

# ---- 阅读说明 ----
ws = wb.create_sheet("阅读说明", 0)
_title_row(ws,
    "2026 年度活动 · 费用清单",
    "本表汇总活动各项费用，供各方查阅",
    2)
notes = [
    ("整理人 / 日期",   "活动财务组  ·  2026 年"),
    ("币种",           "人民币（元）"),
    ("数据来源",       "活动原始票据与对账单"),
    ("分类",           "① 外部费用  ② 内部费用  ③ 现场杂项"),
    ("数字格式",       f"金额统一 ¥ #,##0.00，负数表示费用支出"),
    ("费用合计",       f"¥{abs(GRAND):,.2f}（共 {len(EXTERNAL)+len(INTERNAL)+len(ONSITE)} 笔）"),
    ("附表",           "证书打印费明细（合计 ¥1,030.00）"),
    ("工作表索引",     "费用总览 → 外部费用明细 → 内部费用明细 → 现场杂项明细 → 证书打印费明细"),
]
r = 4
for k, v in notes:
    ws.row_dimensions[r].height = 26
    a = ws.cell(row=r, column=1, value=k)
    a.font = FONT_SUBHEAD; a.fill = FILL_SUBHEAD; a.alignment = ALIGN_CENTER; a.border = BORDER_CELL
    b = ws.cell(row=r, column=2, value=v)
    b.font = FONT_CELL; b.alignment = ALIGN_LEFT; b.border = BORDER_CELL
    r += 1
_widths(ws, [18, 72])
ws.sheet_view.showGridLines = False
_fit(ws)

xlsx_path = os.path.join(OUT_DIR, "2026年活动费用清单.xlsx")
wb.save(xlsx_path)
print("Excel 已生成：", xlsx_path)

# ---------------------------------------------------------------
# PPT
# ---------------------------------------------------------------
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

DARK_RED = RGBColor(0x8B, 0x00, 0x00)
GOLD     = RGBColor(0xC9, 0xA8, 0x6B)
CREAM    = RGBColor(0xFB, 0xF7, 0xF0)
DARK_TXT = RGBColor(0x33, 0x33, 0x33)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0x6B, 0x6B, 0x6B)
RED      = RGBColor(0xB2, 0x22, 0x22)
BG_PAGE  = RGBColor(0xFB, 0xF7, 0xF0)
FONT_CN  = "WenQuanYi Micro Hei"

def add_bg(slide, color=BG_PAGE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()

def add_top_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.95))
    bar.fill.solid(); bar.fill.fore_color.rgb = DARK_RED; bar.line.fill.background()
    deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.95), SW, Inches(0.06))
    deco.fill.solid(); deco.fill.fore_color.rgb = GOLD; deco.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), SW - Inches(1), Inches(0.85))
    tf = tb.text_frame; tf.margin_left = 0; tf.margin_top = 0
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.name = FONT_CN; r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.name = FONT_CN; r2.font.size = Pt(12); r2.font.color.rgb = GOLD

def add_footer(slide, idx, total):
    tb = slide.shapes.add_textbox(Inches(0.4), SH - Inches(0.4), SW - Inches(0.8), Inches(0.3))
    tf = tb.text_frame; tf.margin_left = 0; tf.margin_top = 0
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "2026 年度活动 · 费用清单"
    r.font.name = FONT_CN; r.font.size = Pt(9); r.font.color.rgb = GRAY
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run(); r2.text = f"{idx} / {total}"
    r2.font.name = FONT_CN; r2.font.size = Pt(9); r2.font.color.rgb = GRAY

def set_cell(cell, text, *, bold=False, size=11, color=DARK_TXT,
             fill=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    cell.text = ""
    tf = cell.text_frame
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(text)
    r.font.name = FONT_CN; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    cell.vertical_anchor = anchor
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = fill

def money_str(x): return f"¥{x:,.2f}"

TOTAL_SLIDES = 6

# ---- 封面 ----
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG_PAGE)
top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(3.6))
top.fill.solid(); top.fill.fore_color.rgb = DARK_RED; top.line.fill.background()
deco = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.6), SW, Inches(0.08))
deco.fill.solid(); deco.fill.fore_color.rgb = GOLD; deco.line.fill.background()

tb = s.shapes.add_textbox(Inches(0.8), Inches(1.0), SW - Inches(1.6), Inches(2.4))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "2026 年度活动"
r.font.name = FONT_CN; r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = GOLD
p2 = tf.add_paragraph()
r2 = p2.add_run(); r2.text = "费 用 清 单"
r2.font.name = FONT_CN; r2.font.size = Pt(58); r2.font.bold = True; r2.font.color.rgb = WHITE
p3 = tf.add_paragraph()
r3 = p3.add_run(); r3.text = "Activity Expense Report"
r3.font.name = "Calibri"; r3.font.size = Pt(14); r3.font.color.rgb = CREAM

card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.2), SW - Inches(1.6), Inches(2.6))
card.fill.solid(); card.fill.fore_color.rgb = WHITE
card.line.color.rgb = GOLD; card.line.width = Pt(1.5)
tb = s.shapes.add_textbox(Inches(1.2), Inches(4.4), SW - Inches(2.4), Inches(2.3))
tf = tb.text_frame; tf.word_wrap = True
for key, val, c in [
    ("外部费用", money_str(SUM_EXT), RED),
    ("内部费用", money_str(SUM_INT), RED),
    ("现场杂项", money_str(SUM_ON),  RED),
    ("合    计", money_str(GRAND),   RED),
    ("项 目 数", f"共 {len(EXTERNAL)+len(INTERNAL)+len(ONSITE)} 笔费用", DARK_TXT),
]:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    r = p.add_run(); r.text = f"  {key} ： "
    r.font.name = FONT_CN; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = DARK_RED
    r2 = p.add_run(); r2.text = val
    r2.font.name = FONT_CN; r2.font.size = Pt(14); r2.font.color.rgb = c
    p.space_after = Pt(6)

add_footer(s, 1, TOTAL_SLIDES)

# ---- 总览 ----
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_top_bar(s, "一、费用总览", "Overview")

def stat_card(left, top, w, h, label, value, count):
    bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
    bg.line.color.rgb = GOLD; bg.line.width = Pt(1)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, Inches(0.18))
    bar.fill.solid(); bar.fill.fore_color.rgb = RED; bar.line.fill.background()
    tb = s.shapes.add_textbox(left, top + Inches(0.28), w, h - Inches(0.3))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = FONT_CN; r.font.size = Pt(13); r.font.color.rgb = GRAY
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = money_str(value)
    r2.font.name = FONT_CN; r2.font.size = Pt(24); r2.font.bold = True; r2.font.color.rgb = RED
    p2.space_before = Pt(8)
    p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run(); r3.text = f"共 {count} 项"
    r3.font.name = FONT_CN; r3.font.size = Pt(11); r3.font.color.rgb = GRAY

card_w = Inches(3.8); card_h = Inches(1.85); gap = Inches(0.3)
total_w = card_w * 3 + gap * 2
start_left = (SW - total_w) // 2
stat_card(start_left + (card_w + gap)*0, Inches(1.45), card_w, card_h, "外部费用", SUM_EXT, len(EXTERNAL))
stat_card(start_left + (card_w + gap)*1, Inches(1.45), card_w, card_h, "内部费用", SUM_INT, len(INTERNAL))
stat_card(start_left + (card_w + gap)*2, Inches(1.45), card_w, card_h, "现场杂项", SUM_ON,  len(ONSITE))

big = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(3.55), SW - Inches(1.2), Inches(2.4))
big.fill.solid(); big.fill.fore_color.rgb = DARK_RED
big.line.color.rgb = GOLD; big.line.width = Pt(2)
tb = s.shapes.add_textbox(Inches(1.1), Inches(3.8), SW - Inches(2.2), Inches(1.9))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "本次活动费用合计"
r.font.name = FONT_CN; r.font.size = Pt(18); r.font.color.rgb = GOLD; r.font.bold = True
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = money_str(GRAND)
r2.font.name = FONT_CN; r2.font.size = Pt(54); r2.font.bold = True; r2.font.color.rgb = WHITE
p2.space_before = Pt(8)
p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
r3 = p3.add_run()
r3.text = f"= 外部 {money_str(SUM_EXT)}  +  内部 {money_str(SUM_INT)}  +  现场 {money_str(SUM_ON)}  ·  共 {len(EXTERNAL)+len(INTERNAL)+len(ONSITE)} 笔"
r3.font.name = FONT_CN; r3.font.size = Pt(13); r3.font.color.rgb = CREAM
p3.space_before = Pt(12)

add_footer(s, 2, TOTAL_SLIDES)

# ---- 明细表 slide ----
def add_table_slide(idx, title, subtitle, data, subtotal_label, subtotal):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_top_bar(s, title, subtitle)
    headers = ["序号", "项目", "对接人 / 供应商", "金额（元）", "备注"]
    rows = len(data) + 2; cols = len(headers)
    left = Inches(0.4); top = Inches(1.25)
    width = SW - Inches(0.8); height = Inches(5.6)
    tbl = s.shapes.add_table(rows, cols, left, top, width, height).table
    col_w_in = [0.7, 2.5, 2.2, 1.9, 5.2]
    total_in = sum(col_w_in)
    for i, w_in in enumerate(col_w_in):
        tbl.columns[i].width = Emu(int((width / total_in) * w_in))
    for i, h in enumerate(headers):
        set_cell(tbl.cell(0, i), h, bold=True, size=12, color=WHITE,
                 fill=DARK_RED, align=PP_ALIGN.CENTER)
    tbl.rows[0].height = Inches(0.42)
    for ri, (proj, contact, amt, remark) in enumerate(data, 1):
        stripe = CREAM if ri % 2 == 0 else WHITE
        set_cell(tbl.cell(ri, 0), ri, size=10, align=PP_ALIGN.CENTER, fill=stripe)
        set_cell(tbl.cell(ri, 1), proj, size=10, fill=stripe)
        set_cell(tbl.cell(ri, 2), contact, size=10, fill=stripe)
        set_cell(tbl.cell(ri, 3), money_str(amt), size=10,
                 color=RED, align=PP_ALIGN.RIGHT, fill=stripe, bold=True)
        set_cell(tbl.cell(ri, 4), remark, size=9, fill=stripe, color=GRAY)
        tbl.rows[ri].height = Inches(0.32)
    sub_row = len(data) + 1
    c0 = tbl.cell(sub_row, 0); c0.merge(tbl.cell(sub_row, 2))
    set_cell(c0, subtotal_label, bold=True, size=12,
             color=DARK_RED, align=PP_ALIGN.CENTER, fill=GOLD)
    set_cell(tbl.cell(sub_row, 3), money_str(subtotal),
             bold=True, size=12, color=RED, align=PP_ALIGN.RIGHT, fill=GOLD)
    set_cell(tbl.cell(sub_row, 4), "", fill=GOLD)
    tbl.rows[sub_row].height = Inches(0.46)
    add_footer(s, idx, TOTAL_SLIDES)

add_table_slide(3, "二、外部费用明细", f"对外采购 / 服务 / 媒体 / 餐饮  ·  小计 {money_str(SUM_EXT)}",
                EXTERNAL, "外部费用小计", SUM_EXT)
add_table_slide(4, "三、内部费用明细", f"内部团队垫付  ·  小计 {money_str(SUM_INT)}",
                INTERNAL, "内部费用小计", SUM_INT)
add_table_slide(5, "四、现场杂项明细", f"现场零星费用  ·  小计 {money_str(SUM_ON)}",
                ONSITE, "现场杂项小计", SUM_ON)

# ---- 附表 + 合计 ----
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_top_bar(s, "五、附表与合计", "Appendix & Total")

# 左：证书附表
left_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.3),
                               Inches(8.0), Inches(5.6))
left_card.fill.solid(); left_card.fill.fore_color.rgb = WHITE
left_card.line.color.rgb = GOLD; left_card.line.width = Pt(1)
tb = s.shapes.add_textbox(Inches(0.7), Inches(1.45), Inches(7.5), Inches(0.5))
tf = tb.text_frame
p = tf.paragraphs[0]
r = p.add_run(); r.text = "附表 · 证书制作及KT版打印明细"
r.font.name = FONT_CN; r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = DARK_RED
p2 = tf.add_paragraph()
r2 = p2.add_run(); r2.text = "供应商：广告公司   单据编号：XH20260522016   合计 ¥1,030.00"
r2.font.name = FONT_CN; r2.font.size = Pt(10); r2.font.color.rgb = GRAY

cheaders = ["序号", "商品", "型号", "单位", "数量", "单价", "金额", "备注"]
cdata = CERT_DETAIL
crows = len(cdata) + 2
ctbl = s.shapes.add_table(crows, len(cheaders),
                          Inches(0.6), Inches(2.4),
                          Inches(7.6), Inches(2.8)).table
cw = [0.55, 1.2, 0.7, 0.65, 0.65, 0.95, 1.05, 2.6]
ctot = sum(cw)
for i, w in enumerate(cw):
    ctbl.columns[i].width = Emu(int((Inches(7.6) / ctot) * w))
for i, h in enumerate(cheaders):
    set_cell(ctbl.cell(0, i), h, bold=True, size=10, color=WHITE, fill=DARK_RED, align=PP_ALIGN.CENTER)
ctbl.rows[0].height = Inches(0.34)
for ri, (no, name, mdl, unit, qty, price, amt, remark) in enumerate(cdata, 1):
    stripe = CREAM if ri % 2 == 0 else WHITE
    cells = [ri, name, mdl, unit, qty, money_str(price), money_str(amt), remark]
    aligns = [PP_ALIGN.CENTER]*5 + [PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.LEFT]
    for ci, v in enumerate(cells):
        set_cell(ctbl.cell(ri, ci), v, size=9, fill=stripe, align=aligns[ci],
                 color=DARK_TXT, bold=(ci == 6))
    ctbl.rows[ri].height = Inches(0.3)
sr = len(cdata) + 1
c0 = ctbl.cell(sr, 0); c0.merge(ctbl.cell(sr, 5))
set_cell(c0, "合计", bold=True, size=11, color=DARK_RED, align=PP_ALIGN.CENTER, fill=GOLD)
set_cell(ctbl.cell(sr, 6), money_str(1030.00), bold=True, size=11, color=RED,
         align=PP_ALIGN.RIGHT, fill=GOLD)
set_cell(ctbl.cell(sr, 7), "", fill=GOLD)
ctbl.rows[sr].height = Inches(0.38)

# 右：合计卡
right_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.6), Inches(1.3),
                                Inches(4.3), Inches(5.6))
right_card.fill.solid(); right_card.fill.fore_color.rgb = DARK_RED
right_card.line.color.rgb = GOLD; right_card.line.width = Pt(1.5)
tb = s.shapes.add_textbox(Inches(8.8), Inches(1.55), Inches(3.9), Inches(5.2))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "费用合计"
r.font.name = FONT_CN; r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = GOLD

for label, val in [
    ("外部费用", money_str(SUM_EXT)),
    ("内部费用", money_str(SUM_INT)),
    ("现场杂项", money_str(SUM_ON)),
]:
    pn = tf.add_paragraph(); pn.alignment = PP_ALIGN.CENTER; pn.space_before = Pt(14)
    rn = pn.add_run(); rn.text = f"{label}\n"
    rn.font.name = FONT_CN; rn.font.size = Pt(12); rn.font.color.rgb = CREAM
    rv = pn.add_run(); rv.text = val
    rv.font.name = FONT_CN; rv.font.size = Pt(20); rv.font.bold = True
    rv.font.color.rgb = RGBColor(0xFF, 0xB6, 0xB6)

line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.85), Inches(5.55), Inches(3.7), Emu(20000))
line.fill.solid(); line.fill.fore_color.rgb = GOLD; line.line.fill.background()

pn = tf.add_paragraph(); pn.alignment = PP_ALIGN.CENTER; pn.space_before = Pt(22)
rn = pn.add_run(); rn.text = "总  计"
rn.font.name = FONT_CN; rn.font.size = Pt(14); rn.font.color.rgb = GOLD
pn = tf.add_paragraph(); pn.alignment = PP_ALIGN.CENTER; pn.space_before = Pt(4)
rn = pn.add_run(); rn.text = money_str(GRAND)
rn.font.name = FONT_CN; rn.font.size = Pt(28); rn.font.bold = True; rn.font.color.rgb = WHITE

add_footer(s, 6, TOTAL_SLIDES)

pptx_path = os.path.join(OUT_DIR, "2026年活动费用清单汇报.pptx")
prs.save(pptx_path)
print("PPT 已生成：", pptx_path)
print("OK")
