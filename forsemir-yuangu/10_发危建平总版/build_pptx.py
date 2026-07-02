"""元谷 2 万方招商策略汇报 — 危建平总汇报版 (脱敏外发).

要点:
  - 称谓统一『危建平总』
  - 胡教授团队 = 第三方专业招商运营服务机构 (代表复旦住房政策研究中心 + 上海市科企联), 不提新设公司
  - 租金动态平衡: 招商期均 1.5-1.8 → 长期稳定 2.0-2.5, 保底 2.0 含物业
  - 仲量联行爬楼大数据 = 我方核心优势
  - 服务费用: 取消月费 (不签对赌); 沙龙费单收一次性; 挂牌 10 万/项可选; 超额奖励适度
  - 以招商策略为主, 脱敏, 便于汇报
"""
from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path(__file__).with_name("元谷2万方招商策略汇报(危建平总汇报版).pptx")

PRIMARY = RGBColor(0x0F, 0x24, 0x4E)
PRIMARY2 = RGBColor(0x1B, 0x3A, 0x6B)
ACCENT = RGBColor(0xF2, 0x7E, 0x2D)
GOLD = RGBColor(0xC9, 0xA2, 0x4B)
TEAL = RGBColor(0x18, 0x8B, 0x8B)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x1E, 0x8E, 0x5A)
LIGHT = RGBColor(0xF4, 0xF6, 0xFA)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD8, 0xDF, 0xEA)
TEXT = RGBColor(0x21, 0x2B, 0x42)
MUTED = RGBColor(0x66, 0x70, 0x86)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

_page = {"n": 0}


def _font(run, size, *, bold=False, color=TEXT, name="微软雅黑"):
    run.font.name = name; run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color
    from pptx.oxml.ns import qn
    rpr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea'):
        el = rpr.find(qn(tag))
        if el is None:
            el = rpr.makeelement(qn(tag), {}); rpr.append(el)
        el.set('typeface', name)


def _rect(slide, l, t, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE, line_w=1.0):
    sp = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def _text(slide, l, t, w, h, lines, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, spec in enumerate(lines):
        text, size, bold, color = spec[0], spec[1], spec[2], spec[3]
        sa = spec[4] if len(spec) > 4 else 3
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sa); p.space_before = Pt(0)
        r = p.add_run(); r.text = text; _font(r, size, bold=bold, color=color)
    return tb


def bg(slide, color=LIGHT):
    _rect(slide, -0.1, -0.1, 13.6, 7.7, fill=color)


def header(slide, kicker, title, *, accent=ACCENT):
    _rect(slide, 0, 0, 13.333, 1.15, fill=PRIMARY)
    _rect(slide, 0, 1.15, 13.333, 0.06, fill=accent)
    _rect(slide, 0.55, 0.26, 0.12, 0.62, fill=accent)
    _text(slide, 0.85, 0.20, 11.2, 0.34, [(kicker, 11, False, RGBColor(0xB9, 0xC6, 0xDC))])
    _text(slide, 0.85, 0.46, 11.4, 0.6, [(title, 22, True, WHITE)])
    _text(slide, 10.4, 0.20, 2.4, 0.8, [("元谷 YUANGU", 11.5, True, RGBColor(0xE8, 0xC8, 0x7A)),
                                        ("胡教授团队 · 第三方服务", 8.5, False, RGBColor(0xB9, 0xC6, 0xDC))],
          align=PP_ALIGN.RIGHT)


def footer(slide):
    _page["n"] += 1
    _rect(slide, 0, 7.28, 13.333, 0.22, fill=PRIMARY)
    _text(slide, 0.55, 7.28, 10.0, 0.22, [("元谷项目 2 万方招商策略汇报 · 危建平总汇报版 (脱敏)", 8, False, RGBColor(0xC6, 0xD0, 0xE2))], anchor=MSO_ANCHOR.MIDDLE)
    _text(slide, 11.8, 7.28, 1.0, 0.22, [(f"{_page['n']:02d}", 8, True, RGBColor(0xE8, 0xC8, 0x7A))], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def kpi(slide, l, t, w, h, label, value, sub, *, accent=ACCENT):
    _rect(slide, l, t, w, h, fill=CARD, line=LINE, line_w=0.75)
    _rect(slide, l, t, 0.1, h, fill=accent)
    _text(slide, l + 0.25, t + 0.16, w - 0.4, 0.3, [(label, 10.5, False, MUTED)])
    _text(slide, l + 0.25, t + 0.46, w - 0.4, 0.5, [(value, 20, True, PRIMARY)])
    _text(slide, l + 0.25, t + h - 0.42, w - 0.4, 0.36, [(sub, 9.5, False, MUTED)])


def chip(slide, l, t, w, h, text, *, fill=PRIMARY, fg=WHITE, size=11, bold=True):
    sp = _rect(slide, l, t, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = sp.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; _font(r, size, bold=bold, color=fg)
    return sp


def table(slide, l, t, w, headers, rows, *, col_w=None, header_fill=PRIMARY,
          row_h=0.42, head_h=0.42, fs=10.5, hfs=11, accent_first=False):
    n = len(rows) + 1; cols = len(headers)
    gt = slide.shapes.add_table(n, cols, Inches(l), Inches(t), Inches(w), Inches(head_h + row_h * len(rows))).table
    gt.first_row = False; gt.horz_banding = False
    if col_w:
        for j, cw in enumerate(col_w): gt.columns[j].width = Inches(cw)
    gt.rows[0].height = Inches(head_h)
    for j, htxt in enumerate(headers):
        c = gt.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = header_fill
        c.margin_left = Inches(0.1); c.margin_right = Inches(0.06)
        c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02); c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]; r = p.add_run(); r.text = htxt; _font(r, hfs, bold=True, color=WHITE)
    for i, row in enumerate(rows, start=1):
        gt.rows[i].height = Inches(row_h)
        for j, val in enumerate(row):
            c = gt.cell(i, j); c.fill.solid()
            c.fill.fore_color.rgb = (RGBColor(0xED, 0xF1, 0xF8) if i % 2 == 0 else CARD)
            c.margin_left = Inches(0.1); c.margin_right = Inches(0.06)
            c.margin_top = Inches(0.01); c.margin_bottom = Inches(0.01); c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]; r = p.add_run(); r.text = val
            ia = accent_first and j == 0
            _font(r, fs, bold=ia, color=(PRIMARY if ia else TEXT))
    return gt


def bullets(slide, l, t, w, h, items, *, fs=14, gap=6, color=TEXT, marker="—", mc=ACCENT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.margin_left = 0; tf.margin_top = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.space_before = Pt(0)
        rm = p.add_run(); rm.text = marker + "  "; _font(rm, fs, bold=True, color=mc)
        r = p.add_run(); r.text = it; _font(r, fs, bold=False, color=color)
    return tb


def divider(prs, no, zh, en, pts):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, -0.1, -0.1, 13.6, 7.7, fill=PRIMARY)
    _rect(s, 0, 0, 4.6, 7.5, fill=PRIMARY2)
    _rect(s, 0.9, 2.0, 0.9, 0.9, fill=ACCENT, shape=MSO_SHAPE.OVAL)
    _text(s, 0.9, 2.0, 0.9, 0.9, [(no, 26, True, WHITE)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _text(s, 0.9, 3.15, 3.4, 1.6, [(zh, 28, True, WHITE), (en, 11.5, False, RGBColor(0x9F, 0xB2, 0xD0), 0)])
    yy = 2.1
    for pt in pts:
        _rect(s, 5.4, yy + 0.12, 0.16, 0.16, fill=ACCENT, shape=MSO_SHAPE.OVAL)
        _text(s, 5.8, yy, 6.9, 0.6, [(pt, 14, False, RGBColor(0xE5, 0xEB, 0xF5))])
        yy += 0.74
    return s


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ===== P1 封面 =====
    s = prs.slides.add_slide(blank)
    _rect(s, -0.1, -0.1, 13.6, 7.7, fill=PRIMARY)
    _rect(s, 0, 0, 13.333, 0.18, fill=ACCENT)
    _rect(s, 0, 7.32, 13.333, 0.18, fill=GOLD)
    _rect(s, 9.3, -0.1, 4.3, 7.7, fill=PRIMARY2)
    _rect(s, 11.0, 1.2, 1.7, 1.7, fill=ACCENT, shape=MSO_SHAPE.OVAL)
    _rect(s, 10.2, 4.4, 1.0, 1.0, fill=GOLD, shape=MSO_SHAPE.OVAL)
    _rect(s, 11.9, 5.3, 0.7, 0.7, fill=TEAL, shape=MSO_SHAPE.OVAL)
    _rect(s, 0.9, 2.3, 0.16, 2.1, fill=ACCENT)
    _text(s, 1.25, 2.2, 8.0, 0.5, [("SEMIR 森马 (上海) 国际运营中心 · 元谷项目", 14, False, RGBColor(0xB9, 0xC6, 0xDC))])
    _text(s, 1.2, 2.65, 8.2, 2.2, [
        ("元谷 2 万方招商策略汇报", 36, True, WHITE),
        ("产业研发办公 · 招商运营服务方案", 20, True, RGBColor(0xE8, 0xC8, 0x7A), 6),
    ])
    _text(s, 1.25, 4.75, 8.0, 0.5, [("IP + AI 双轨  ·  动态租金平衡  ·  仲量联行爬楼大数据赋能", 14.5, False, RGBColor(0xCF, 0xD8, 0xE8))])
    chip(s, 1.25, 5.6, 5.4, 0.55, "胡教授团队(代表复旦住房政策研究中心 · 上海市科企联)", fill=ACCENT, size=11)
    _text(s, 1.25, 6.4, 8.0, 0.5, [("呈:危建平总     ·     2026 年     ·     本汇报仅用于商务沟通", 11.5, False, RGBColor(0x9F, 0xB2, 0xD0))])

    # ===== P2 议程 =====
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "AGENDA", "汇报议程")
    ag = [
        ("01", "项目与机会", "双重身份 / 2 万方产业研发"),
        ("02", "我方服务能力", "第三方专业机构 / 仲量联行大数据优势"),
        ("03", "招商战略", "IP+AI 双轨 / 四级漏斗 / 六大出彩点"),
        ("04", "租金与节点", "动态租金平衡 / 9·30 与 5·1 硬节点"),
        ("05", "品牌与活动", "5 项挂牌 / 6 场沙龙 / 5·22 峰会"),
        ("06", "服务费用与价值", "佣金 / 沙龙 / 挂牌 / 对项目方贡献"),
    ]
    x0, y0, cw, ch, gx, gy = 0.7, 1.55, 6.05, 1.42, 0.25, 0.22
    for i, (no, zh, en) in enumerate(ag):
        col = i % 2; rowi = i // 2
        l = x0 + col * (cw + gx); t = y0 + rowi * (ch + gy)
        _rect(s, l, t, cw, ch, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, l, t, 1.1, ch, fill=PRIMARY)
        _text(s, l, t, 1.1, ch, [(no, 24, True, RGBColor(0xE8, 0xC8, 0x7A))], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, l + 1.3, t + 0.32, cw - 1.45, 0.5, [(zh, 16, True, PRIMARY)])
        _text(s, l + 1.3, t + 0.82, cw - 1.45, 0.45, [(en, 10.5, False, MUTED)])
    footer(s)

    # ===== 01 =====
    divider(prs, "01", "项目与机会", "PROJECT & OPPORTUNITY", [
        "元谷:大零号湾文创融合核心区 + 上海唯一科技时尚特色小镇",
        "本期聚焦 4#+5# 楼 5F+ 共约 2 万方产业研发办公",
        "15 号线元江路站 TOD, 单日客流 5-7 万人次",
        "危建平总已锁定保底租金 2.0 元/㎡/天(含物业)",
    ])

    s = prs.slides.add_slide(blank); bg(s)
    header(s, "01 · 项目与机会", "元谷的双重战略身份与产业底盘")
    kpi(s, 0.7, 1.5, 2.95, 1.5, "总建筑面积", "22 万㎡", "1#-6# 共 6 栋")
    kpi(s, 3.83, 1.5, 2.95, 1.5, "本期招商范围", "2 万㎡", "4#+5# 楼 5F+", accent=GOLD)
    kpi(s, 6.96, 1.5, 2.95, 1.5, "TOD 单日客流", "5-7 万", "15 号线元江路站", accent=TEAL)
    kpi(s, 10.09, 1.5, 2.55, 1.5, "15min 覆盖", "24万+12万", "居住 + 产业人口", accent=PRIMARY2)
    chip(s, 0.7, 3.35, 5.9, 0.62, "大零号湾文创融合核心区(闵行五大中心之一)", fill=PRIMARY, size=12.5)
    chip(s, 6.75, 3.35, 5.9, 0.62, "上海市唯一科技时尚特色小镇", fill=ACCENT, size=12.5)
    bullets(s, 0.7, 4.3, 12.0, 2.4, [
        "能级比肩漕河泾、张江;享区级最高 1,000 万 + 市级最高 5,000 万产业专项",
        "4# 楼 5F+ 定位 AI 主轴(国际创意层);5# 楼 5F+ 定位 IP 主轴(产业总部层)",
        "产业愿景:打造 AI + IP 双赛道协同的产业研发园区,与森马时尚科技定位契合",
        "危建平总已锁定保底租金 2.0 元/㎡/天(含物业费)作为长期稳定锚点",
    ], fs=14)
    footer(s)

    # ===== 02 =====
    divider(prs, "02", "我方服务能力", "OUR CAPABILITY", [
        "第三方专业招商运营服务机构",
        "代表复旦大学住房政策研究中心 + 上海市科技企业联合会",
        "★ 仲量联行爬楼大数据 — 我方核心招商优势",
        "腾讯算力 / 产业基金 / 福布斯 / 北欧会客厅 资源",
    ])

    s = prs.slides.add_slide(blank); bg(s)
    header(s, "02 · 我方服务能力", "胡教授团队:第三方专业招商运营服务机构")
    _rect(s, 0.7, 1.5, 11.95, 1.25, fill=PRIMARY)
    _text(s, 1.0, 1.5, 11.4, 1.25, [
        ("胡教授团队 · 代表复旦大学住房政策研究中心、上海市科技企业联合会", 15, True, WHITE),
        ("以第三方专业服务身份, 为元谷项目提供产业招商策略 + 政府政策对接 + 品牌活动 + 资源导入的全链条服务。", 11.5, False, RGBColor(0xCF, 0xD8, 0xE8), 0),
    ], anchor=MSO_ANCHOR.MIDDLE)
    res = [
        ("★ 仲量联行爬楼大数据", "200+ 家精准客户清单, 转化率 +30% (我方独家优势)", ACCENT),
        ("腾讯算力补贴", "新注册 AI 公司 3 月合同免费 / 算力 85 折", TEAL),
        ("产业基金 (追觅 + CVC)", "返投落地, 资本驱动签约", GOLD),
        ("北欧创新国际会客厅", "海外 IP / 技术进入中国首站", PRIMARY2),
        ("福布斯产业影响力奖", "国际品牌势能 + 年度评选 IP", ACCENT),
        ("复旦 + 上海市科企联", "学术背书 + 政府对接 + 政策代办", TEAL),
    ]
    x0, y0, cw, ch, gx, gy = 0.7, 3.05, 3.85, 1.55, 0.18, 0.2
    for i, (t1, t2, ac) in enumerate(res):
        col = i % 3; rowi = i // 3
        l = x0 + col * (cw + gx); t = y0 + rowi * (ch + gy)
        _rect(s, l, t, cw, ch, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, l, t, cw, 0.1, fill=ac)
        _text(s, l + 0.22, t + 0.26, cw - 0.4, 0.5, [(t1, 13, True, PRIMARY)])
        _text(s, l + 0.22, t + 0.78, cw - 0.4, 0.7, [(t2, 10.5, False, MUTED)])
    footer(s)

    # P 仲量联行优势专页
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "02 · 我方服务能力", "核心优势:仲量联行爬楼大数据赋能精准招商")
    kpi(s, 0.7, 1.5, 3.85, 1.5, "目标客户清单", "200+ 家", "精准画像, 直接触达", accent=ACCENT)
    kpi(s, 4.75, 1.5, 3.85, 1.5, "转化率提升", "+30%", "VS 普通广撒网招商", accent=GREEN)
    kpi(s, 8.8, 1.5, 3.85, 1.5, "覆盖区域", "闵行+临港+漕河泾+张江", "存量企业全画像", accent=TEAL)
    bullets(s, 0.7, 3.35, 12, 3.2, [
        "爬楼大数据 = 对上海存量产业楼宇的企业逐层扫描, 形成『谁在哪栋楼、面积多大、何时到期、付租能力』的精准画像;",
        "我方据此锁定真正有迁址/扩租需求、且付租能力强的中大型 AI / IP 企业, 把行业 110:1 的转化率压缩到更高效区间;",
        "针对元谷 2 万方, 优先匹配 2,000-20,000㎡ 大客户(1 家大户 ≫ 10 家小户), 提升招商速度与稳定性;",
        "该数据能力为我方独家优势, 直接赋能元谷招商, 是普通招商团队不具备的弹药库。",
    ], fs=13.5, gap=10)
    footer(s)

    # ===== 03 =====
    divider(prs, "03", "招商战略", "LEASING STRATEGY", [
        "战略主轴:IP + AI 双轨 (AI 拿政策/资本, IP 兼顾文化)",
        "客户策略:优先 2,000-20,000㎡ 大客户(转化率 110:1)",
        "四级招商漏斗 + 招商六大『出彩点』",
        "五档客户配比, 与森马原产业规划一致",
    ])

    s = prs.slides.add_slide(blank); bg(s)
    header(s, "03 · 招商战略", "战略主轴:IP + AI 双轨")
    table(s, 0.7, 1.5, 11.95,
          ["维度", "纯二次元(不推荐)", "IP + AI 双轨(推荐)"],
          [
              ["招商面", "200-500㎡ 小户, 招满需 400-800 家", "AI/IP 中大型, 招满仅需 30-60 家"],
              ["支付能力", "单家年租 < 30 万", "单家年租 100-500 万(真金白银)"],
              ["政府政策", "拿政策难, 创新基金不投", "享高新 + 专精特新 + 创新券 + 闵行专项"],
              ["与森马协同", "二次元与商业难协同", "IP+AI 与森马时尚科技定位契合"],
              ["先例", "无 2 万方成功案例", "杨浦 1 万方『AI+IP』已通过先例"],
          ],
          col_w=[1.8, 4.85, 5.3], row_h=0.6, accent_first=True)
    chip(s, 0.7, 5.0, 11.95, 0.7, "结论:主轴改 AI 拿政府资源 + 资本, 保留 IP 概念(视频/漫画/直播/动漫全链条)兼顾文化导向", fill=ACCENT, size=13)
    bullets(s, 0.7, 5.95, 12, 1.1, [
        "4# 楼 5F+ = AI 主轴(国际创意层);5# 楼 5F+ = IP 主轴(产业总部层)",
        "行业转化率 110:1 → 招 1 家大户 ≫ 招 10 家小户, 效率差 5 倍",
    ], fs=13)
    footer(s)

    s = prs.slides.add_slide(blank); bg(s)
    header(s, "03 · 招商战略", "四级招商漏斗")
    funnel = [
        ("L1 牌照锚定", "头部央企 / 行业协会", "AI 潮玩产业基地 + 潮玩次元专委 双牌照", "6,000㎡", ACCENT),
        ("L2 资本招商", "中型 AI / IP 企业", "产业基金 + 腾讯算力 + 返投落地", "5,000㎡", GOLD),
        ("L3 大数据爬楼", "小型潮玩 / 服务机构", "★ 仲量联行爬楼大数据(我方优势)", "6,000-8,000㎡", TEAL),
        ("L4 活动带流", "中小型服务机构", "6 场沙龙 + 5·22 峰会 + 福布斯", "3,000-4,000㎡", PRIMARY2),
    ]
    yy = 1.55
    for name, who, how, area, ac in funnel:
        _rect(s, 0.7, yy, 11.95, 1.12, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, 0.7, yy, 2.5, 1.12, fill=ac)
        _text(s, 0.7, yy, 2.5, 1.12, [(name, 15, True, WHITE)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, 3.4, yy + 0.16, 6.6, 0.4, [(who, 13, True, PRIMARY)])
        _text(s, 3.4, yy + 0.58, 6.8, 0.45, [(how, 11, False, MUTED)])
        chip(s, 10.4, yy + 0.33, 2.0, 0.46, area, fill=ac, size=12)
        yy += 1.24
    _text(s, 0.7, 6.6, 12, 0.5, [("漏斗合计可产生 20,000-23,000㎡ 招商管道 → 实际签约 18,000-20,000㎡(满租)", 12, True, ACCENT)])
    footer(s)

    s = prs.slides.add_slide(blank); bg(s)
    header(s, "03 · 招商战略", "招商六大『出彩点』(VS 普通招商)")
    pts = [
        ("① 牌照即招商", "双牌照前置, 客户『送进来』不是『拉进来』"),
        ("② 基金即招商", "产业基金返投 + 腾讯算力, 资本驱动签约"),
        ("③ 数据即招商", "★ 仲量联行爬楼大数据, 转化率 +30%"),
        ("④ 峰会即招商", "5·22 峰会 200+ VIP, 1 天锁 30-50 客户"),
        ("⑤ 沙龙即招商", "6 场 × ≥30 客户 = 5-12 家直接成果"),
        ("⑥ 学术即招商", "复旦 + 上海交大 + 上海市科企联 背书"),
    ]
    x0, y0, cw, ch, gx, gy = 0.7, 1.55, 5.9, 1.42, 0.15, 0.2
    for i, (t1, t2) in enumerate(pts):
        col = i % 2; rowi = i // 2
        l = x0 + col * (cw + gx); t = y0 + rowi * (ch + gy)
        _rect(s, l, t, cw, ch, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, l, t, 0.1, ch, fill=ACCENT)
        _text(s, l + 0.28, t + 0.2, cw - 0.5, 0.5, [(t1, 15, True, PRIMARY)])
        _text(s, l + 0.28, t + 0.74, cw - 0.5, 0.55, [(t2, 11.5, False, MUTED)])
    _text(s, 0.7, 6.95, 12, 0.4, [("普通园区只做 ③④;我方六维齐发 → 招商速度领先 6-12 个月, 单方租金溢价更快兑现", 12, True, ACCENT)])
    footer(s)

    # ===== 04 =====
    divider(prs, "04", "租金与节点", "RENT & MILESTONES", [
        "动态租金平衡:招商期 1.5-1.8 → 长期稳定 2.0-2.5",
        "保底 2.0 元/㎡/天(含物业, 危建平总已签)",
        "硬节点一:2026/9/30 完成 2,000㎡ 签约",
        "硬节点二:2027/5/1 项目开业, 50%+ 签约",
    ])

    # P 动态租金
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "04 · 租金与节点", "动态租金平衡模型:先低价招满, 后抬升")
    kpi(s, 0.7, 1.5, 3.85, 1.6, "招商期(满租前)", "1.5-1.8 元", "先低价快速招满", accent=ACCENT)
    kpi(s, 4.75, 1.5, 3.85, 1.6, "长期稳定(满租后)", "2.0-2.5 元", "逐步抬升至稳定区间", accent=GREEN)
    kpi(s, 8.8, 1.5, 3.85, 1.6, "保底锚点(已签)", "2.0 元含物业", "危建平总锁定的长期数", accent=PRIMARY2)
    _text(s, 0.7, 3.4, 12, 0.4, [("策略逻辑(单位:元/㎡/天)", 14, True, PRIMARY)])
    # 简易阶梯条形示意
    stages = [("招商初期\n0-9 月", 1.5, ACCENT), ("加速期\n9-18 月", 1.7, ACCENT), ("满租稳定\n18-24 月", 2.1, GREEN), ("长期\n24 月+", 2.4, GREEN)]
    bx = 0.9; bw = 2.7; maxv = 2.6
    for label, v, ac in stages:
        bh = (v / maxv) * 2.1
        _rect(s, bx, 3.95 + (2.1 - bh), bw - 0.4, bh, fill=ac)
        _text(s, bx, 3.95 + (2.1 - bh) - 0.32, bw - 0.4, 0.3, [(f"{v} 元", 13, True, PRIMARY)], align=PP_ALIGN.CENTER)
        _text(s, bx, 6.15, bw - 0.4, 0.6, [(label, 10.5, False, MUTED)], align=PP_ALIGN.CENTER)
        bx += bw
    _text(s, 0.7, 6.85, 12, 0.4, [("先以 1.5-1.8 元低价快速招满形成产业氛围, 再借满租势能与产业生态逐步抬升至 2.0-2.5 元, 长期锚定保底 2 元(含物业)", 11, True, ACCENT)])
    footer(s)

    # P 两节点 + 五阶段
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "04 · 租金与节点", "两个硬节点 · 倒推全部排期")
    kpi(s, 0.7, 1.5, 5.9, 1.4, "硬节点一 · 2026/9/30 (T+100 天)", "2,000㎡ 签约", "可含直播基地 / 共享设计中心", accent=RED)
    kpi(s, 6.75, 1.5, 5.9, 1.4, "硬节点二 · 2027/5/1 (T+314 天)", "开业 + 50% 签约", "2 万方达成 50% + 消费氛围", accent=ACCENT)
    table(s, 0.7, 3.15, 11.95,
          ["阶段", "时间窗口", "累计签约目标", "里程碑"],
          [
              ["阶段 0 准备", "启动 - 第 3 周", "0", "团队就位 + 资源接入 + 物料"],
              ["阶段 1 抢节点", "第 3 周 - 9/30", "2,300㎡(含擦边球)", "★ 9/30 硬节点达成"],
              ["阶段 2 加速", "10/1 - 12/31", "5,500㎡(27%)", "爬楼 + 沙龙 + 挂牌"],
              ["阶段 3 开业", "1/1 - 5/1", "10,400㎡(52%)", "★ 5/1 开业典礼"],
              ["阶段 4 满租", "5/2 - 次年", "18,000㎡+(90%)", "沙龙 IP 化 + 抬租"],
          ],
          col_w=[2.0, 2.6, 3.35, 4.0], row_h=0.5, accent_first=True)
    footer(s)

    # P 月度签约率
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "04 · 租金与节点", "月度签约率推进(招商期按 1.7 元估算)")
    table(s, 0.7, 1.5, 11.95,
          ["月份", "本月新增", "累计签约", "签约率", "里程碑"],
          [
              ["第 1 月", "300㎡", "300㎡", "1.5%", "战队就位 + 物料"],
              ["第 2 月", "800㎡", "1,100㎡", "5.5%", "L1 牌照 + 大客户接触"],
              ["第 3 月 (9/30)", "1,200㎡", "2,300㎡", "11.5%", "★ 硬节点一达成"],
              ["第 6 月", "—", "5,500㎡", "27.5%", "Q4 攻坚收官"],
              ["第 10 月 (5/1)", "1,500㎡", "10,400㎡", "52%", "★ 硬节点二达成"],
              ["第 11 月", "1,500㎡", "11,900㎡", "59.5%", "开业典礼 + 启动抬租"],
          ],
          col_w=[2.4, 2.0, 2.3, 1.85, 3.4], row_h=0.5, accent_first=True)
    chip(s, 0.7, 5.3, 11.95, 0.66, "9/30 节点 2,000㎡ 来源:大客户 1,000-2,000㎡ + 直播基地 800-1,000㎡ + AI 设计中心 500-800㎡ + 服务中心 300-500㎡ → 中性期望 3,200-5,800㎡, 概率 >95%", fill=PRIMARY, size=11)
    _text(s, 0.7, 6.25, 12, 0.5, [("满租后启动抬租, 把招商期 1.5-1.8 元逐步推向长期稳定 2.0-2.5 元, 为项目方持续增收", 12.5, True, GREEN)])
    footer(s)

    # ===== 05 =====
    divider(prs, "05", "品牌与活动", "BRAND & EVENTS", [
        "5 项产业牌照(可选, 10 万/项)",
        "6 场产业沙龙, 每场 ≥30 个目标产业客户",
        "5·22 AI 商业化峰会借势, 200+ VIP 转化",
        "活动 IP 化 → 年度品牌势能 + 媒体声量",
    ])

    s = prs.slides.add_slide(blank); bg(s)
    header(s, "05 · 品牌与活动", "5 项产业牌照(可选 · 10 万/项)")
    bands = [
        ("①", "AI 潮玩产业基地", "中国动漫集团"),
        ("②", "潮玩次元商业专委会", "中国百货商业协会"),
        ("③", "复旦大学住房政策研究中心 · 元谷分中心", "复旦大学"),
        ("④", "上海市科技企业联合会 · 元谷产业基地", "上海市科企联"),
        ("⑤", "福布斯产业影响力奖 · 元谷专场", "福布斯"),
    ]
    yy = 1.55
    for no, name, body in bands:
        _rect(s, 0.7, yy, 11.95, 0.86, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, 0.7, yy, 0.86, 0.86, fill=PRIMARY)
        _text(s, 0.7, yy, 0.86, 0.86, [(no, 22, True, RGBColor(0xE8, 0xC8, 0x7A))], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, 1.75, yy + 0.16, 8.0, 0.55, [(name, 14, True, PRIMARY)])
        _text(s, 1.75, yy + 0.5, 8.0, 0.32, [("出牌方:" + body, 10.5, False, MUTED)])
        chip(s, 10.7, yy + 0.24, 1.8, 0.38, "10 万 · 可选", fill=ACCENT, size=11)
        yy += 0.98
    _text(s, 0.7, 6.5, 12, 0.5, [("挂牌费 10 万元/项, 由项目方按需选择, 选定后挂牌前一次性付清(最多 5 项)。", 12, True, PRIMARY)])
    footer(s)

    s = prs.slides.add_slide(blank); bg(s)
    header(s, "05 · 品牌与活动", "6 场产业沙龙(每场 ≥30 客户) + 5·22 峰会")
    table(s, 0.7, 1.5, 7.55,
          ["#", "时间", "主题", "联办方"],
          [
              ["1", "第 1 月", "AI+潮玩(借势 5·22)", "中动漫+腾讯"],
              ["2", "第 3 月", "潮玩出海", "北欧+福布斯"],
              ["3", "第 5 月", "投融资路演", "产业基金+银行"],
              ["4", "第 7 月", "设计与创意", "上海交大+科企联"],
              ["5", "第 9 月", "内容IP·Z世代", "中百协+中动漫"],
              ["6", "第 11 月", "政策补贴·小镇", "闵行科协+复旦"],
          ],
          col_w=[0.5, 1.4, 2.65, 3.0], row_h=0.5, fs=10, accent_first=True)
    _rect(s, 8.55, 1.5, 4.1, 4.5, fill=PRIMARY)
    _rect(s, 8.55, 1.5, 4.1, 0.6, fill=ACCENT)
    _text(s, 8.75, 1.5, 3.8, 0.6, [("5·22 AI 商业化峰会借势", 13, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, 8.78, 2.3, 3.7, 3.6, [
        "200+ VIP(北大/复旦/银行/投资机构)",
        "元谷设专属展位 + 闭门招商台",
        "增设产业影响力榜 + 福布斯背书",
        "峰会嘉宾 → 沙龙 #1 主题嘉宾",
        "沙龙费由我方单独收取(不分润)",
    ], fs=11.5, gap=10, color=RGBColor(0xE5, 0xEB, 0xF5), marker="▸", mc=RGBColor(0xE8, 0xC8, 0x7A))
    _text(s, 0.7, 6.3, 12, 0.5, [("每场沙龙触达 ≥30 家产业客户 + ≥100 万次媒体曝光;6 场打包沙龙执行费一次性付清。", 12, True, ACCENT)])
    footer(s)

    # ===== 05b 媒体宣传 (专页) =====
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "05 · 品牌与活动", "媒体宣传:为招商蓄势、为品牌立势")
    _text(s, 0.7, 1.35, 12, 0.4, [("『主流媒体权威引领 + 社交媒体精准触达』双轮驱动, 放大声量、赋能招商", 13, True, PRIMARY)])
    # 左:主流媒体
    _rect(s, 0.7, 1.95, 5.85, 4.4, fill=CARD, line=LINE, line_w=0.75)
    _rect(s, 0.7, 1.95, 5.85, 0.55, fill=PRIMARY)
    _text(s, 0.95, 1.95, 5.4, 0.55, [("① 中央级 / 全国性 / 上海主流媒体", 13, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, 0.95, 2.7, 5.4, 3.5, [
        "中央级:央视新闻、央广网、人民网、新华网、中新网、中国日报",
        "全国财经政经:第一财经、21 世纪经济报道、澎湃、界面、上证报、中证报",
        "上海区域:解放日报、文汇报、上观新闻、闵行区政府官网",
        "作用:权威背书、定调舆情、提升招商公信力",
    ], fs=12, gap=12)
    # 右:社交媒体
    _rect(s, 6.75, 1.95, 5.9, 4.4, fill=CARD, line=LINE, line_w=0.75)
    _rect(s, 6.75, 1.95, 5.9, 0.55, fill=ACCENT)
    _text(s, 7.0, 1.95, 5.4, 0.55, [("② 社交媒体流量投放(双轮驱动)", 13, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, 7.0, 2.7, 5.5, 3.5, [
        "抖音:高频互动 + 信息流广告, 精准算法推荐",
        "今日头条:深度阅读场景, 广度覆盖",
        "投放区域:上海;精准锁定目标受众",
        "为招商转化奠定流量基础与认知铺垫",
    ], fs=12, gap=12)
    _text(s, 0.7, 6.55, 12, 0.5, [("重大节点(9/30 签约、5/1 开业)前集中投放, 把品牌声量直接转化为招商线索;合作结束提供完整投放报告(链接+曝光+分析)", 11, True, ACCENT)])
    footer(s)

    # 媒体服务内容与报价
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "05 · 品牌与活动", "媒体宣传服务内容与报价(可选模块)")
    table(s, 0.7, 1.55, 11.95,
          ["服务项目", "数量", "费用 / 说明"],
          [
              ["原创内容生产(品牌+营销稿原创采写)", "3 篇", "折后 5 万元(媒体总曝光量预计 ≥ 150 万)"],
              ["中央级/全国/上海主流媒体宣发", "10 篇", "打包报价(按媒体档位据实核算)"],
              ["社交媒体流量投放(抖音/今日头条信息流, 上海)", "按投放量", "按投放量计(短视频拍摄制作费不含)"],
              ["营销策划与舆情管理咨询", "不超 2 次", "含于媒体服务包"],
          ],
          col_w=[5.6, 2.0, 4.35], row_h=0.7, accent_first=True)
    bullets(s, 0.7, 5.0, 12, 1.5, [
        "服务周期:自签订起至 2026/12/31, 乙方享优先续约权;",
        "该媒体宣传为可选模块, 与招商佣金、沙龙、挂牌相互独立、不影响其他计费;",
        "联动 6 场沙龙 + 5·22 峰会 + 5 项挂牌 + 福布斯, 形成『活动造势 + 媒体放大』闭环。",
    ], fs=12.5, gap=8)
    footer(s)

    # ===== 06 服务费用与价值 =====
    divider(prs, "06", "服务费用与价值", "FEES & VALUE", [
        "已取消月费(不签对赌), 让项目方更轻装",
        "招商佣金:实际成交年租金的 1.5-2 个月",
        "沙龙费一次性付清, 挂牌 10 万/项可选",
        "对项目方:满租年租金 + 资产增值",
    ])

    s = prs.slides.add_slide(blank); bg(s)
    header(s, "06 · 服务费用与价值", "服务费用构成(已取消月费 · 不签对赌)")
    table(s, 0.7, 1.5, 11.95,
          ["费用类别", "标准", "支付方式"],
          [
              ["招商佣金(核心)", "实际成交年租金的 1.5 / 1.75 / 2.0 个月(按面积阶梯)", "起租后结算"],
              ["沙龙执行费", "6 场打包 30 万(我方单独收取, 不与他方分润)", "★ 一次性付清(或分两次)"],
              ["挂牌费(可选)", "10 万元/项 × 项目方选定项数(最多 5 项)", "★ 挂牌前一次性付清"],
              ["媒体宣传(可选模块)", "原创 3 篇 5 万(曝光≥150万)+ 主流媒体宣发 + 社交投放", "按服务包结算(详见媒体方案)"],
              ["超额奖励(适度)", "满租率 ≥ 90% 后给予, 金额适度", "正常分期支付(不要求一次性)"],
          ],
          col_w=[2.6, 6.7, 2.65], row_h=0.56, accent_first=True)
    bullets(s, 0.7, 4.55, 12, 2.0, [
        "取消固定月费 → 我方不要求对赌、不绑保底, 项目方零固定支出压力, 完全按招商成果付费;",
        "沙龙执行费由我方单独收取、不与他方分润, 建议一次性付清以减少后续财务往来;",
        "挂牌费 10 万/项由项目方按需选择, 灵活可控;超额奖励仅在满租达标后适度计提、正常给付。",
    ], fs=13, gap=9)
    footer(s)

    s = prs.slides.add_slide(blank); bg(s)
    header(s, "06 · 服务费用与价值", "对项目方的价值:动态租金下的持续增收")
    kpi(s, 0.7, 1.55, 3.85, 1.6, "满租期年租金", "约 1,460-1,825 万", "2 万㎡ × 2.0-2.5 元", accent=GREEN)
    kpi(s, 4.75, 1.55, 3.85, 1.6, "保底锚点", "≈ 1,460 万/年", "2.0 元含物业 × 2 万㎡", accent=PRIMARY2)
    kpi(s, 8.8, 1.55, 3.85, 1.6, "资产增值", "数千万级", "产业认证 + TOD 板块溢价", accent=ACCENT)
    table(s, 0.7, 3.5, 11.95,
          ["阶段", "租金水平", "年化租金(2 万㎡)", "说明"],
          [
              ["招商期(先低招满)", "1.5-1.8 元/㎡/天", "约 1,095-1,314 万", "快速形成产业氛围"],
              ["满租稳定(后抬升)", "2.0-2.5 元/㎡/天", "约 1,460-1,825 万", "动态平衡, 持续增收"],
              ["长期保底锚点", "2.0 元含物业", "约 1,460 万", "危建平总已签的稳定数"],
          ],
          col_w=[3.2, 2.8, 3.15, 2.8], row_h=0.52, accent_first=True)
    _text(s, 0.7, 6.35, 12, 0.5, [("我方完全按招商成果收费, 与项目方利益高度绑定:招得越快越满, 项目方租金兑现越早越稳。", 12.5, True, GREEN)])
    footer(s)

    # P 协议要点 + 下一步
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "07 · 合作要点与下一步", "合作协议要点 + 启动建议")
    _text(s, 0.7, 1.45, 12, 0.4, [("合作协议核心要点(详见附件协议)", 14, True, PRIMARY)])
    bullets(s, 0.7, 1.95, 12, 2.6, [
        "服务方:胡教授团队(代表复旦大学住房政策研究中心、上海市科技企业联合会), 以第三方专业服务身份提供招商运营;",
        "服务范围:元谷 4#+5# 楼 5F+ 共约 2 万㎡产业研发办公的招商策划、执行、政府对接、品牌活动;",
        "★ 排他性:在协议期内, 项目方授予我方该范围独家招商运营权, 不另行委托第三方机构;",
        "服务费用:招商佣金(核心)+ 沙龙费(一次性)+ 挂牌费(10 万/项可选)+ 超额奖励(适度), 无月费;",
        "租金策略:执行动态租金平衡(先低招满、后抬升), 长期锚定保底 2.0 元(含物业)。",
    ], fs=12.5, gap=8)
    _text(s, 0.7, 5.0, 12, 0.4, [("启动建议", 14, True, PRIMARY)])
    bullets(s, 0.7, 5.45, 12, 1.5, [
        "建议尽快确定合作意向并签署协议, 启动后第 3 周完成团队就位 + 资源接入;",
        "赶 2026/9/30 硬节点(2,000㎡ 签约), 错过 7 月窗口节点风险陡增;",
        "我方可先提供招商物料与策略, 与项目方并行推进合同流程。",
    ], fs=12.5, gap=8)
    footer(s)

    # 尾页
    s = prs.slides.add_slide(blank)
    _rect(s, -0.1, -0.1, 13.6, 7.7, fill=PRIMARY)
    _rect(s, 0, 0, 13.333, 0.18, fill=ACCENT)
    _rect(s, 0, 7.32, 13.333, 0.18, fill=GOLD)
    _rect(s, 1.2, 2.4, 0.16, 2.2, fill=ACCENT)
    _text(s, 1.55, 2.5, 11.0, 2.6, [
        ("先低价招满 · 后动态抬租", 38, True, WHITE),
        ("IP + AI 双轨 · 仲量联行大数据赋能 · 专业第三方服务", 18, False, RGBColor(0xE8, 0xC8, 0x7A), 10),
        ("以招商成果说话, 与项目方利益高度绑定", 15, False, RGBColor(0xCF, 0xD8, 0xE8), 0),
    ])
    chip(s, 1.55, 5.65, 6.2, 0.55, "胡教授团队(代表复旦住房政策研究中心 · 上海市科企联)", fill=ACCENT, size=12)
    _text(s, 1.55, 6.5, 11, 0.4, [("呈:危建平总  ·  期待与项目方达成合作", 12, False, RGBColor(0x9F, 0xB2, 0xD0))])

    prs.save(OUT)
    print(f"Wrote {OUT} ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
