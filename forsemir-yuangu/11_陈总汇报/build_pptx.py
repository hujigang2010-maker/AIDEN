"""元谷项目 · 致森马资管陈总汇报 PPT.

目标:回应陈总对『合作必要性』的疑虑,讲清四件事:
  1. 与森马合作的必要性
  2. 政府政策资源端优势
  3. 沙龙带来的优质客群
  4. 对项目本身的宣传作用

语境:已向森马董事长汇报过;对方知悉我方为上海科协/科企联体系团队。
"""
from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path(__file__).with_name("元谷_致陈总_合作必要性汇报.pptx")

PRIMARY = RGBColor(0x0F, 0x24, 0x4E)
PRIMARY2 = RGBColor(0x1B, 0x3A, 0x6B)
ACCENT = RGBColor(0xF2, 0x7E, 0x2D)
GOLD = RGBColor(0xC9, 0xA2, 0x4B)
TEAL = RGBColor(0x18, 0x8B, 0x8B)
GREEN = RGBColor(0x1E, 0x8E, 0x5A)
RED = RGBColor(0xC0, 0x39, 0x2B)
LIGHT = RGBColor(0xF4, 0xF6, 0xFA)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD8, 0xDF, 0xEA)
TEXT = RGBColor(0x21, 0x2B, 0x42)
MUTED = RGBColor(0x66, 0x70, 0x86)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

_page = {"n": 0}


def _font(run, size, *, bold=False, color=TEXT, name="微软雅黑"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    from pptx.oxml.ns import qn
    rpr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea"):
        el = rpr.find(qn(tag))
        if el is None:
            el = rpr.makeelement(qn(tag), {})
            rpr.append(el)
        el.set("typeface", name)


def _rect(slide, l, t, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE, line_w=1.0):
    sp = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def _text(slide, l, t, w, h, lines, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, spec in enumerate(lines):
        text, size, bold, color = spec[0], spec[1], spec[2], spec[3]
        sa = spec[4] if len(spec) > 4 else 3
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(sa)
        p.space_before = Pt(0)
        r = p.add_run()
        r.text = text
        _font(r, size, bold=bold, color=color)
    return tb


def bg(slide, color=LIGHT):
    _rect(slide, -0.1, -0.1, 13.6, 7.7, fill=color)


def header(slide, kicker, title, *, accent=ACCENT):
    _rect(slide, 0, 0, 13.333, 1.15, fill=PRIMARY)
    _rect(slide, 0, 1.15, 13.333, 0.06, fill=accent)
    _rect(slide, 0.55, 0.26, 0.12, 0.62, fill=accent)
    _text(slide, 0.85, 0.20, 11.2, 0.34, [(kicker, 11, False, RGBColor(0xB9, 0xC6, 0xDC))])
    _text(slide, 0.85, 0.46, 11.4, 0.6, [(title, 22, True, WHITE)])
    _text(
        slide,
        10.2,
        0.20,
        2.6,
        0.8,
        [
            ("元谷 YUANGU", 11.5, True, RGBColor(0xE8, 0xC8, 0x7A)),
            ("致陈总 · 合作必要性", 8.5, False, RGBColor(0xB9, 0xC6, 0xDC)),
        ],
        align=PP_ALIGN.RIGHT,
    )


def footer(slide):
    _page["n"] += 1
    _rect(slide, 0, 7.28, 13.333, 0.22, fill=PRIMARY)
    _text(
        slide,
        0.55,
        7.28,
        10.0,
        0.22,
        [("元谷项目 · 致森马资管陈总 · 合作必要性专场汇报", 8, False, RGBColor(0xC6, 0xD0, 0xE2))],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    _text(
        slide,
        11.8,
        7.28,
        1.0,
        0.22,
        [(f"{_page['n']:02d}", 8, True, RGBColor(0xE8, 0xC8, 0x7A))],
        align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def kpi(slide, l, t, w, h, label, value, sub, *, accent=ACCENT):
    _rect(slide, l, t, w, h, fill=CARD, line=LINE, line_w=0.75)
    _rect(slide, l, t, 0.1, h, fill=accent)
    _text(slide, l + 0.25, t + 0.16, w - 0.4, 0.3, [(label, 10.5, False, MUTED)])
    _text(slide, l + 0.25, t + 0.46, w - 0.4, 0.5, [(value, 20, True, PRIMARY)])
    _text(slide, l + 0.25, t + 1.0, w - 0.4, 0.35, [(sub, 10.5, False, MUTED)])


def chip(slide, l, t, w, h, text, *, fill=ACCENT, size=12):
    _rect(slide, l, t, w, h, fill=fill)
    _text(slide, l, t, w, h, [(text, size, True, WHITE)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def bullets(slide, l, t, w, h, items, *, fs=13, gap=8, color=TEXT, marker="●", mc=ACCENT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.space_before = Pt(0)
        r0 = p.add_run()
        r0.text = f"{marker} "
        _font(r0, fs, bold=True, color=mc)
        r1 = p.add_run()
        r1.text = item
        _font(r1, fs, bold=False, color=color)
    return tb


def table(slide, l, t, w, headers, rows, *, col_w=None, row_h=0.55, fs=11, accent_first=False):
    n_cols = len(headers)
    n_rows = 1 + len(rows)
    if col_w is None:
        col_w = [w / n_cols] * n_cols
    shape = slide.shapes.add_table(n_rows, n_cols, Inches(l), Inches(t), Inches(w), Inches(row_h * n_rows))
    tbl = shape.table
    for i, cw in enumerate(col_w):
        tbl.columns[i].width = Inches(cw)
    for j, htxt in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = htxt
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                _font(r, fs, bold=True, color=WHITE)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if j == 0 else PP_ALIGN.LEFT
                for r in p.runs:
                    bold = accent_first and j == 0
                    _font(r, fs, bold=bold, color=PRIMARY if bold else TEXT)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xEF, 0xF2, 0xF8) if i % 2 else CARD
    return shape


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    _page["n"] = 0

    # ===== 封面 =====
    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, 13.333, 7.5, fill=PRIMARY)
    _rect(s, 0, 0, 0.35, 7.5, fill=ACCENT)
    _text(s, 1.1, 1.5, 11, 0.4, [("森马（上海）国际运营中心 · 元谷项目", 14, False, RGBColor(0xB9, 0xC6, 0xDC))])
    _text(s, 1.1, 2.1, 11, 1.0, [("为什么需要我们一起做?", 36, True, WHITE)])
    _text(
        s,
        1.1,
        3.2,
        11,
        0.5,
        [("合作必要性 · 政府政策资源 · 沙龙优质客群 · 项目宣传作用", 16, False, RGBColor(0xE8, 0xC8, 0x7A))],
    )
    _rect(s, 1.1, 4.0, 2.2, 0.08, fill=ACCENT)
    _text(
        s,
        1.1,
        4.4,
        10,
        1.0,
        [
            ("致：森马资管 陈总", 16, True, WHITE),
            ("汇报方：胡教授团队（复旦大学住房政策研究中心 · 上海市科技企业联合会）", 13, False, RGBColor(0xCF, 0xD8, 0xE8)),
            ("背景：已向森马董事长汇报；贵方知悉我方为上海科协体系团队", 12, False, RGBColor(0x9F, 0xB2, 0xD0)),
        ],
    )
    _text(s, 1.1, 6.6, 10, 0.4, [("仅用于商务沟通 · 第三方专业服务视角", 11, False, RGBColor(0x7A, 0x8C, 0xA8))])

    # ===== 目录 =====
    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "AGENDA", "今天只回答四个问题")
    qs = [
        ("01", "合作必要性", "为什么不是自己招、而是必须一起做?", ACCENT),
        ("02", "政府政策资源", "科协/科企联体系能给项目什么真金白银?", TEAL),
        ("03", "沙龙优质客群", "活动不是热闹, 是可转化的客群管道", GOLD),
        ("04", "项目宣传作用", "我们如何把元谷做成『看得见的品牌』?", PRIMARY2),
    ]
    for i, (no, title, sub, ac) in enumerate(qs):
        y = 1.55 + i * 1.25
        _rect(s, 0.7, y, 11.95, 1.1, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, 0.7, y, 1.4, 1.1, fill=ac)
        _text(s, 0.7, y, 1.4, 1.1, [(no, 26, True, WHITE)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, 2.4, y + 0.22, 9.5, 0.4, [(title, 18, True, PRIMARY)])
        _text(s, 2.4, y + 0.62, 9.5, 0.35, [(sub, 13, False, MUTED)])
    footer(s)

    # ===== 01 必要性总论 =====
    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "01 · 合作必要性", "一句话：元谷要的不是『多一个人招租』，而是『少走 6–12 个月弯路』")
    cols = [
        ("自己做", "常见路径", RED, [
            "广撒网中介 + 熟人介绍",
            "行业转化率约 110:1",
            "政策申报靠企业自己摸",
            "活动热闹, 难沉淀客户",
            "品牌声量靠零散发稿",
        ]),
        ("与我们合作", "加速路径", GREEN, [
            "精准客群 + 政府通道并行",
            "把 110:1 压到可运营区间",
            "高新/专精特新等代办陪跑",
            "沙龙即招商, 单场 ≥30 客户",
            "活动造势 + 媒体放大闭环",
        ]),
    ]
    for i, (title, sub, ac, items) in enumerate(cols):
        x = 0.7 + i * 6.2
        _rect(s, x, 1.5, 5.9, 5.4, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, x, 1.5, 5.9, 0.85, fill=ac)
        _text(s, x + 0.3, 1.55, 5.3, 0.45, [(title, 18, True, WHITE)])
        _text(s, x + 0.3, 1.95, 5.3, 0.35, [(sub, 12, False, RGBColor(0xFF, 0xF0, 0xE8))])
        bullets(s, x + 0.35, 2.6, 5.3, 4.0, items, fs=14, gap=14, mc=ac)
    footer(s)

    # ===== 110:1 漏斗 =====
    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "01 · 合作必要性", "核心证据：行业 110:1 客群漏斗 —— 自己招，成本在『看不见的地方』")
    kpi(s, 0.7, 1.45, 3.85, 1.55, "行业转化基准", "110 : 1", "触达 110 家 ≈ 签 1 家", accent=RED)
    kpi(s, 4.75, 1.45, 3.85, 1.55, "大户策略效率差", "× 5 倍", "1 家大户 ≫ 10 家小户", accent=ACCENT)
    kpi(s, 8.8, 1.45, 3.85, 1.55, "我方压缩路径", "数据 + 政策 + 沙龙", "把漏斗从『碰运气』变『可运营』", accent=GREEN)

    steps = [
        ("触达层", "110+", "广撒网 / 中介名单 / 陌生拜访", MUTED),
        ("意向层", "20-30", "到访、看房、要方案", PRIMARY2),
        ("谈判层", "5-8", "比租金、比政策、比配套", TEAL),
        ("签约层", "1", "真正落定的成功客群", ACCENT),
    ]
    _text(s, 0.7, 3.25, 12, 0.35, [("漏斗示意（产业园区招商常见结构）", 13, True, PRIMARY)])
    x = 0.7
    widths = [3.4, 2.9, 2.5, 2.2]
    for (name, num, desc, ac), ww in zip(steps, widths):
        _rect(s, x, 3.7, ww, 2.2, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, x, 3.7, ww, 0.12, fill=ac)
        _text(s, x + 0.15, 3.95, ww - 0.3, 0.35, [(name, 12, False, MUTED)], align=PP_ALIGN.CENTER)
        _text(s, x + 0.15, 4.35, ww - 0.3, 0.55, [(num, 28, True, PRIMARY)], align=PP_ALIGN.CENTER)
        _text(s, x + 0.15, 5.1, ww - 0.3, 0.6, [(desc, 11, False, MUTED)], align=PP_ALIGN.CENTER)
        x += ww + 0.15
    _text(
        s,
        0.7,
        6.15,
        12,
        0.7,
        [
            ("结论：2 万方若走小户路线，相当于要跑数百上千次触达；合作价值不在『多招几家』，而在『用更短路径筛出付租能力强的成功客群』。", 13, True, ACCENT),
        ],
    )
    footer(s)

    # ===== 必要性四维 =====
    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "01 · 合作必要性", "我们补的是森马资管『自有招商』最难补的四块短板")
    cards = [
        ("政府通道", "科协 / 科企联 / 闵行专项", "高新、专精特新、创新券、小镇政策 —— 企业愿意为『能办成事』而来", TEAL),
        ("精准客群", "仲量联行爬楼大数据", "200+ 家画像清单，转化率约 +30%，把 110:1 往下压", ACCENT),
        ("活动获客", "6 场产业沙龙 + 峰会", "单场 ≥30 目标客户，活动不是热闹，是签约管道", GOLD),
        ("品牌立势", "挂牌 + 媒体放大", "权威背书 + 曝光蓄势，让元谷在开业前就『被看见』", PRIMARY2),
    ]
    for i, (t1, t2, t3, ac) in enumerate(cards):
        col, rowi = i % 2, i // 2
        l = 0.7 + col * 6.2
        t = 1.5 + rowi * 2.55
        _rect(s, l, t, 5.95, 2.35, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, l, t, 5.95, 0.12, fill=ac)
        _text(s, l + 0.3, t + 0.35, 5.4, 0.4, [(t1, 18, True, PRIMARY)])
        chip(s, l + 0.3, t + 0.9, 4.6, 0.4, t2, fill=ac, size=11)
        _text(s, l + 0.3, t + 1.5, 5.4, 0.65, [(t3, 13, False, MUTED)])
    footer(s)

    # ===== 02 政府资源总览 =====
    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "02 · 政府政策资源", "我们不是『介绍政策』，而是『能对接、能申报、能落地』")
    _text(
        s,
        0.7,
        1.4,
        12,
        0.4,
        [("身份前提：上海科协体系团队 · 上海市科技企业联合会 · 复旦大学住房政策研究中心（董事长汇报时已说明）", 13, True, PRIMARY)],
    )
    items = [
        ("上海市科企联", "产业组织通道", "元谷产业基地挂牌候选；科技企业资源池与活动联办", TEAL),
        ("闵行 / 区级专项", "属地政策落地", "对接科委、商务、街道；特色小镇与文创扶持线索", ACCENT),
        ("复旦学术平台", "公信力背书", "住房政策研究中心 · 元谷分中心；智库活动与课题联动", GOLD),
        ("申报陪跑能力", "企业真金白银", "高新、专精特新、创新券等 —— 客户因政策而来、因落地而留", PRIMARY2),
    ]
    for i, (a, b, c, ac) in enumerate(items):
        y = 2.0 + i * 1.15
        _rect(s, 0.7, y, 11.95, 1.0, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, 0.7, y, 0.14, 1.0, fill=ac)
        _text(s, 1.1, y + 0.15, 3.2, 0.7, [(a, 15, True, PRIMARY)], anchor=MSO_ANCHOR.MIDDLE)
        chip(s, 4.4, y + 0.28, 2.4, 0.44, b, fill=ac, size=11)
        _text(s, 7.0, y + 0.2, 5.3, 0.65, [(c, 13, False, MUTED)], anchor=MSO_ANCHOR.MIDDLE)
    footer(s)

    # ===== 政策申报明细 =====
    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "02 · 政府政策资源", "政策申报端：把『政策故事』变成客户决策理由")
    table(
        s,
        0.7,
        1.45,
        11.95,
        ["政策类型", "对企业的价值", "对元谷招商的作用"],
        [
            ["高新技术企业认定", "税率优惠 / 资质背书", "吸引有成长诉求的科技企业入驻"],
            ["专精特新 / 小巨人", "专项支持与品牌势能", "锁定优质中型产业客户"],
            ["创新券 / 研发补贴", "直接降本", "缩短看房到签约的决策周期"],
            ["闵行区文创 / 小镇专项", "属地扶持与落地配套", "强化大零号湾区位政策叙事"],
            ["算力 / 产业基金返投", "资本与算力驱动", "AI 主轴客户『带着资源来』"],
        ],
        col_w=[3.4, 3.8, 4.75],
        row_h=0.72,
        accent_first=True,
    )
    chip(
        s,
        0.7,
        5.6,
        11.95,
        0.7,
        "对陈总：政策不是附属服务，而是提升转化率的『硬弹药』—— 同一套房，谁能帮客户办成事，谁就先签下",
        fill=PRIMARY,
        size=13,
    )
    _text(
        s,
        0.7,
        6.5,
        12,
        0.4,
        [("配合 IP+AI 双轨：AI 侧吃透政策与资本，IP 侧兼顾文化叙事，避免纯二次元难拿政策的困局。", 12.5, True, ACCENT)],
    )
    footer(s)

    # ===== 03 沙龙客群 =====
    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "03 · 沙龙优质客群", "沙龙不是『办活动』，而是可计量的优质客群入口")
    kpi(s, 0.7, 1.45, 3.85, 1.5, "年度沙龙", "6 场", "主题招商，不是泛活动", accent=ACCENT)
    kpi(s, 4.75, 1.45, 3.85, 1.5, "单场目标客户", "≥ 30 家", "产业相关、可跟进", accent=TEAL)
    kpi(s, 8.8, 1.45, 3.85, 1.5, "年度直接成果预期", "5–12 家", "可进入谈判的成果客户", accent=GREEN)

    table(
        s,
        0.7,
        3.2,
        11.95,
        ["场次", "主题", "联办方", "客群画像"],
        [
            ["#1", "AI + 潮玩（借势峰会）", "中动漫 + 腾讯", "AI/内容科技中大型"],
            ["#2", "潮玩出海", "北欧会客厅 + 福布斯", "出海 IP / 外贸型"],
            ["#3", "投融资路演", "产业基金 + 银行", "有融资诉求的成长企业"],
            ["#4", "设计与创意", "交大 + 科企联", "设计总部 / 创意机构"],
            ["#5", "内容 IP · Z 世代", "中百协 + 中动漫", "内容、直播、IP 运营"],
            ["#6", "政策补贴 · 小镇", "闵行科协 + 复旦", "政策敏感型科技企业"],
        ],
        col_w=[1.0, 3.4, 3.5, 4.05],
        row_h=0.48,
        fs=11,
        accent_first=True,
    )
    footer(s)

    # ===== 沙龙与漏斗 =====
    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "03 · 沙龙优质客群", "把 110:1 接到沙龙上：每一场都在『预筛选』成功客群")
    funnel = [
        ("场前邀约", "定向名单", "科企联资源池 + 爬楼清单 + 峰会 VIP", ACCENT),
        ("场中触达", "≥30 家/场", "主题匹配，天然过滤无效流量", TEAL),
        ("场后跟进", "意向池", "到访元谷、政策对接、方案比选", GOLD),
        ("签约转化", "5–12 家/年", "进入正式租赁谈判与落定", GREEN),
    ]
    for i, (a, b, c, ac) in enumerate(funnel):
        x = 0.7 + i * 3.15
        _rect(s, x, 1.55, 3.0, 3.3, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, x, 1.55, 3.0, 0.55, fill=ac)
        _text(s, x, 1.55, 3.0, 0.55, [(a, 14, True, WHITE)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, x + 0.15, 2.4, 2.7, 0.6, [(b, 20, True, PRIMARY)], align=PP_ALIGN.CENTER)
        _text(s, x + 0.2, 3.3, 2.6, 1.2, [(c, 13, False, MUTED)], align=PP_ALIGN.CENTER)
        if i < 3:
            _text(s, x + 2.85, 2.8, 0.4, 0.4, [("→", 20, True, ACCENT)], align=PP_ALIGN.CENTER)
    bullets(
        s,
        0.7,
        5.15,
        12,
        1.6,
        [
            "另借 5·22 AI 商业化峰会：200+ VIP 一天集中触达，嘉宾可转入沙龙 #1；",
            "单场媒体曝光预期 ≥100 万次 —— 客群获客与品牌声量同步发生；",
            "对资管侧：沙龙费用可单独结算，成果可追踪（名单、跟进、到访、签约）。",
        ],
        fs=13.5,
        gap=10,
    )
    footer(s)

    # ===== 04 宣传作用 =====
    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "04 · 项目宣传作用", "我们帮元谷做的，不只是『发稿』，而是『招商前的立势』")
    left = [
        ("权威背书层", "中央级 / 财经 / 上海主流媒体定调", "央视新闻、人民网、新华网、第一财经、解放日报、上观、闵行区政府网等"),
        ("流量触达层", "抖音 + 今日头条信息流", "上海区域精准投放，为招商铺认知"),
    ]
    for i, (t1, t2, t3) in enumerate(left):
        y = 1.45 + i * 2.35
        _rect(s, 0.7, y, 7.4, 2.15, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, 0.7, y, 0.14, 2.15, fill=ACCENT if i == 0 else TEAL)
        _text(s, 1.1, y + 0.25, 6.7, 0.4, [(t1, 16, True, PRIMARY)])
        _text(s, 1.1, y + 0.75, 6.7, 0.35, [(t2, 13, True, ACCENT if i == 0 else TEAL)])
        _text(s, 1.1, y + 1.2, 6.7, 0.7, [(t3, 12.5, False, MUTED)])

    _rect(s, 8.35, 1.45, 4.3, 4.7, fill=PRIMARY)
    _text(s, 8.6, 1.7, 3.9, 0.4, [("宣传对招商的价值", 14, True, RGBColor(0xE8, 0xC8, 0x7A))])
    bullets(
        s,
        8.6,
        2.3,
        3.85,
        3.5,
        [
            "开业前先建立『可信项目』认知",
            "重大节点（9/30、5/1）集中爆发",
            "挂牌 + 沙龙 + 媒体形成闭环",
            "声量可转化为看房线索",
            "降低客户对『新盘』的顾虑",
        ],
        fs=12.5,
        gap=12,
        color=RGBColor(0xE5, 0xEB, 0xF5),
        marker="▸",
        mc=RGBColor(0xE8, 0xC8, 0x7A),
    )
    footer(s)

    # ===== 宣传闭环与可选包 =====
    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "04 · 项目宣传作用", "宣传闭环：活动造势 × 媒体放大 × 招商转化")
    ring = [
        ("① 挂牌立势", "科企联 / 复旦 / 产业基地等牌照", "公信力"),
        ("② 沙龙造势", "6 场产业主题活动", "客群"),
        ("③ 媒体放大", "主流媒体 + 社交投放", "声量"),
        ("④ 招商转化", "名单跟进 → 到访 → 签约", "结果"),
    ]
    for i, (a, b, c) in enumerate(ring):
        x = 0.7 + i * 3.15
        _rect(s, x, 1.5, 3.0, 2.4, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, x, 1.5, 3.0, 0.5, fill=ACCENT if i % 2 == 0 else TEAL)
        _text(s, x, 1.5, 3.0, 0.5, [(a, 13, True, WHITE)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, x + 0.15, 2.25, 2.7, 0.9, [(b, 13, False, TEXT)], align=PP_ALIGN.CENTER)
        chip(s, x + 0.55, 3.35, 1.9, 0.35, c, fill=PRIMARY, size=11)

    table(
        s,
        0.7,
        4.2,
        11.95,
        ["可选媒体模块（示意）", "量级", "说明"],
        [
            ["原创内容（品牌/营销稿）", "3 篇", "折后约 5 万；媒体总曝光预期 ≥150 万"],
            ["主流媒体宣发", "约 10 篇", "中央级 / 全国 / 上海媒体档位据实"],
            ["社交信息流投放", "按量", "抖音 / 头条，上海区域精准"],
        ],
        col_w=[4.5, 2.2, 5.25],
        row_h=0.55,
        accent_first=True,
    )
    _text(
        s,
        0.7,
        6.35,
        12,
        0.4,
        [("说明：媒体为可选模块，可与沙龙、挂牌独立决策；核心仍是用宣传服务招商结果，而不是为热闹买单。", 12.5, True, ACCENT)],
    )
    footer(s)

    # ===== 总收束 =====
    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "结论", "请陈总今天带走的四个判断")
    conclusions = [
        ("必要性", "不是替代资管，而是补齐政策、客群、立势三条腿；否则 110:1 的成本由项目自己扛。", ACCENT),
        ("政策资源", "科协/科企联体系 + 申报陪跑，是客户『愿不愿意来』的关键砝码。", TEAL),
        ("沙龙客群", "6 场 × ≥30 家 = 可追踪的优质管道，预期 5–12 家直接成果。", GOLD),
        ("宣传作用", "挂牌 + 沙龙 + 媒体闭环，让元谷在开业前就被市场看见、被客户信任。", PRIMARY2),
    ]
    for i, (t1, t2, ac) in enumerate(conclusions):
        y = 1.45 + i * 1.2
        _rect(s, 0.7, y, 11.95, 1.05, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, 0.7, y, 2.2, 1.05, fill=ac)
        _text(s, 0.7, y, 2.2, 1.05, [(t1, 16, True, WHITE)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, 3.2, y + 0.25, 9.1, 0.6, [(t2, 14, False, TEXT)], anchor=MSO_ANCHOR.MIDDLE)
    footer(s)

    # ===== 结束页 =====
    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, 13.333, 7.5, fill=PRIMARY)
    _rect(s, 0, 0, 0.35, 7.5, fill=ACCENT)
    _text(s, 1.1, 2.2, 11, 0.6, [("下一步建议", 16, False, RGBColor(0xE8, 0xC8, 0x7A))])
    _text(s, 1.1, 2.9, 11, 1.0, [("先对齐『合作必要性』，再谈执行细节", 28, True, WHITE)])
    bullets(
        s,
        1.1,
        4.2,
        10,
        2.0,
        [
            "若陈总认可四条逻辑：可进入沙龙排期、政策对接清单、首批客群名单对齐；",
            "费用结构已按『轻负担』设计（无月费、无对赌），细节可另场展开；",
            "感谢陈总时间。我们准备好把元谷做成『招得进、留得住、叫得响』的项目。",
        ],
        fs=14,
        gap=12,
        color=RGBColor(0xE5, 0xEB, 0xF5),
        marker="▸",
        mc=RGBColor(0xE8, 0xC8, 0x7A),
    )

    prs.save(OUT)
    print(f"Wrote {OUT} ({_page['n']} content pages + cover/end)")


if __name__ == "__main__":
    build()
