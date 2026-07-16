"""元谷项目 · 致森马资管陈总 · 正式汇报版 PPT.

依据《陈总 9:30 沟通版》PDF 叙事结构系统展开:
  今天只讲一件事 —— 合作必要性与资源价值闭环
  不是“办几场活动”，而是把普通招商升级为
  政策型 × 圈层型 × 传播型招商。

正式汇报口径:
  - 面向森马资管陈总外发
  - 慎用“保证/包拿”措辞
  - 强调可验证试点（首场沙龙 + 首批复访）
"""
from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path(__file__).with_name("元谷_致陈总_合作必要性_正式汇报版.pptx")

PRIMARY = RGBColor(0x0F, 0x24, 0x4E)
PRIMARY2 = RGBColor(0x1B, 0x3A, 0x6B)
ACCENT = RGBColor(0xF2, 0x7E, 0x2D)
GOLD = RGBColor(0xC9, 0xA2, 0x4B)
TEAL = RGBColor(0x18, 0x8B, 0x8B)
GREEN = RGBColor(0x1E, 0x8E, 0x5A)
RED = RGBColor(0xC0, 0x39, 0x2B)
LIGHT = RGBColor(0xF4, 0xF6, 0xFA)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD8, 0xDF, 0xEA)
TEXT = RGBColor(0x21, 0x2B, 0x42)
MUTED = RGBColor(0x66, 0x70, 0x86)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

_page = {"n": 0}


def _font(run, size, *, bold=False, color=TEXT, name="微软雅黑"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    from pptx.oxml.ns import qn
    rpr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea"):
        el = rpr.find(qn(tag))
        if el is None:
            el = rpr.makeelement(qn(tag), {})
            rpr.append(el)
        el.set("typeface", name)


def _rect(slide, l, t, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE, line_w=1.0):
    sp = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def _text(slide, l, t, w, h, lines, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    for i, spec in enumerate(lines):
        text, size, bold, color = spec[0], spec[1], spec[2], spec[3]
        sa = spec[4] if len(spec) > 4 else 3
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(sa)
        p.space_before = Pt(0)
        r = p.add_run()
        r.text = text
        _font(r, size, bold=bold, color=color)
    return tb


def bg(slide, color=LIGHT):
    _rect(slide, -0.1, -0.1, 13.6, 7.7, fill=color)


def header(slide, kicker, title, *, accent=ACCENT):
    _rect(slide, 0, 0, 13.333, 1.15, fill=PRIMARY)
    _rect(slide, 0, 1.15, 13.333, 0.06, fill=accent)
    _rect(slide, 0.55, 0.26, 0.12, 0.62, fill=accent)
    _text(slide, 0.85, 0.20, 10.5, 0.34, [(kicker, 11, False, RGBColor(0xB9, 0xC6, 0xDC))])
    _text(slide, 0.85, 0.46, 10.8, 0.55, [(title, 20, True, WHITE)])
    _text(
        slide,
        10.0,
        0.22,
        2.9,
        0.75,
        [
            ("元谷 YUANGU", 12, True, RGBColor(0xE8, 0xC8, 0x7A)),
            ("正式汇报版 · 致陈总", 9, False, RGBColor(0xB9, 0xC6, 0xDC)),
        ],
        align=PP_ALIGN.RIGHT,
    )


def footer(slide):
    _page["n"] += 1
    _rect(slide, 0, 7.28, 13.333, 0.22, fill=PRIMARY)
    _text(
        slide,
        0.55,
        7.28,
        10.2,
        0.22,
        [("森马（上海）国际运营中心 · 元谷项目 · 合作必要性正式汇报 · 仅供商务沟通", 8, False, RGBColor(0xC6, 0xD0, 0xE2))],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    _text(
        slide,
        11.7,
        7.28,
        1.1,
        0.22,
        [(f"{_page['n']:02d}", 8, True, RGBColor(0xE8, 0xC8, 0x7A))],
        align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def kpi(slide, l, t, w, h, label, value, sub, *, accent=ACCENT):
    _rect(slide, l, t, w, h, fill=CARD, line=LINE, line_w=0.75)
    _rect(slide, l, t, 0.1, h, fill=accent)
    _text(slide, l + 0.22, t + 0.14, w - 0.35, 0.28, [(label, 10.5, False, MUTED)])
    _text(slide, l + 0.22, t + 0.42, w - 0.35, 0.48, [(value, 22, True, PRIMARY)])
    _text(slide, l + 0.22, t + 0.98, w - 0.35, 0.35, [(sub, 10.5, False, MUTED)])


def chip(slide, l, t, w, h, text, *, fill=ACCENT, size=12, color=WHITE):
    _rect(slide, l, t, w, h, fill=fill)
    _text(slide, l, t, w, h, [(text, size, True, color)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def bullets(slide, l, t, w, h, items, *, fs=13, gap=8, color=TEXT, marker="●", mc=ACCENT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.space_before = Pt(0)
        r0 = p.add_run()
        r0.text = f"{marker} "
        _font(r0, fs, bold=True, color=mc)
        r1 = p.add_run()
        r1.text = item
        _font(r1, fs, bold=False, color=color)
    return tb


def table(slide, l, t, w, headers, rows, *, col_w=None, row_h=0.55, fs=11, accent_first=False):
    n_cols = len(headers)
    n_rows = 1 + len(rows)
    if col_w is None:
        col_w = [w / n_cols] * n_cols
    shape = slide.shapes.add_table(n_rows, n_cols, Inches(l), Inches(t), Inches(w), Inches(row_h * n_rows))
    tbl = shape.table
    for i, cw in enumerate(col_w):
        tbl.columns[i].width = Inches(cw)
    for j, htxt in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = htxt
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                _font(r, fs, bold=True, color=WHITE)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if j == 0 else PP_ALIGN.LEFT
                for r in p.runs:
                    bold = accent_first and j == 0
                    _font(r, fs, bold=bold, color=PRIMARY if bold else TEXT)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xEF, 0xF2, 0xF8) if i % 2 else CARD
    return shape


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    _page["n"] = 0

    # ========== 01 封面 ==========
    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, 13.333, 7.5, fill=PRIMARY)
    _rect(s, 0, 0, 0.28, 7.5, fill=ACCENT)
    _rect(s, 0, 6.85, 13.333, 0.65, fill=RGBColor(0x0A, 0x18, 0x36))
    chip(s, 1.1, 1.15, 2.6, 0.42, "陈总沟通 · 正式汇报版", fill=ACCENT, size=12)
    _text(s, 1.1, 1.85, 11, 0.4, [("今天只讲一件事", 16, False, RGBColor(0xE8, 0xC8, 0x7A))])
    _text(s, 1.1, 2.35, 11.2, 1.1, [("元谷项目合作必要性", 36, True, WHITE), ("与资源价值闭环", 36, True, WHITE)])
    _text(
        s,
        1.1,
        4.0,
        11,
        0.9,
        [
            ("不是“办几场活动”，而是用外部资源", 15, False, RGBColor(0xCF, 0xD8, 0xE8)),
            ("把森马项目从普通招商，升级为政策型、圈层型、传播型招商。", 15, False, RGBColor(0xCF, 0xD8, 0xE8)),
        ],
    )
    _text(
        s,
        1.1,
        5.2,
        11,
        0.8,
        [
            ("政府政策资源  ×  沙龙优质客群  ×  项目宣传势能  ×  招商转化漏斗", 13, True, RGBColor(0xE8, 0xC8, 0x7A)),
            ("目标：降低获客成本 · 提高客户信任 · 加快签约转化", 12, False, RGBColor(0x9F, 0xB2, 0xD0)),
        ],
    )
    _text(
        s,
        1.1,
        6.95,
        11,
        0.4,
        [("胡教授团队｜复旦大学住房政策研究中心 / 上海市科技企业联合会相关资源  ·  面向森马资管陈总 · 元谷约 2 万方招商运营合作", 10, False, RGBColor(0x9F, 0xB2, 0xD0))],
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # ========== 02 四个疑虑 ==========
    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "01 · 沟通议程", "先回答陈总的四个疑虑")
    chip(s, 0.7, 1.4, 8.8, 0.42, "把合作从“费用支出”解释成“招商基础设施投入”", fill=PRIMARY, size=13)
    qs = [
        ("01", "为什么要合作？", "森马具备物业与品牌基础，但外部政策通道、产业圈层和第三方背书，需要专业团队补齐。", ACCENT),
        ("02", "政策资源在哪里？", "围绕科创、AI、设计、文创、总部、研发中心、人才、算力、租金补贴等，形成申报与导入路径。", TEAL),
        ("03", "沙龙能带来什么客户？", "不是泛流量，而是企业家、投资机构、政府产业部门、银行与高校校友组成的高信任客群。", GOLD),
        ("04", "宣传有什么作用？", "把元谷从“出租空间”包装成“上海科技时尚产业场景”，降低客户认知成本。", PRIMARY2),
    ]
    for i, (no, title, body, ac) in enumerate(qs):
        col, rowi = i % 2, i // 2
        l = 0.7 + col * 6.2
        t = 2.05 + rowi * 2.25
        _rect(s, l, t, 5.95, 2.05, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, l, t, 5.95, 0.12, fill=ac)
        _text(s, l + 0.25, t + 0.28, 1.0, 0.4, [(no, 18, True, ac)])
        _text(s, l + 1.2, t + 0.32, 4.4, 0.4, [(title, 16, True, PRIMARY)])
        _text(s, l + 0.25, t + 0.9, 5.4, 0.95, [(body, 12.5, False, MUTED)])
    _text(
        s,
        0.7,
        6.55,
        12,
        0.45,
        [("一句话：我们提供的是“招商加速器”，不是单点中介；能把政府资源、精准客群、宣传声量和成交路径合成一个闭环。", 12.5, True, ACCENT)],
    )
    footer(s)

    # ========== 03 合作必要性 ==========
    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "02 · 合作必要性", "空间端强，但招商转化需要“外部信用 + 外部流量 + 外部政策”")
    cols = [
        ("森马已有", "项目硬资产", "元谷约 2 万方空间、森马品牌、TOD 区位与科技时尚定位——这是招商的基础盘。", PRIMARY),
        ("合作补足", "外部资源包", "政府政策申报、科协/科企联圈层、沙龙客群、第三方学术与产业背书。", ACCENT),
        ("共同结果", "招商转化闭环", "从“看场地”变成“看机会”：政策、资源、传播共同降低签约阻力。", GREEN),
    ]
    for i, (a, b, c, ac) in enumerate(cols):
        x = 0.7 + i * 4.15
        _rect(s, x, 1.5, 3.95, 3.6, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, x, 1.5, 3.95, 1.05, fill=ac)
        _text(s, x + 0.2, 1.6, 3.55, 0.4, [(a, 12, False, RGBColor(0xFF, 0xF0, 0xE0))], align=PP_ALIGN.CENTER)
        _text(s, x + 0.2, 1.95, 3.55, 0.5, [(b, 18, True, WHITE)], align=PP_ALIGN.CENTER)
        _text(s, x + 0.25, 2.85, 3.45, 2.0, [(c, 14, False, TEXT)])
    _rect(s, 0.7, 5.35, 11.95, 1.5, fill=PRIMARY)
    _text(s, 1.0, 5.5, 11.4, 0.35, [("关键判断", 12, True, RGBColor(0xE8, 0xC8, 0x7A))])
    _text(
        s,
        1.0,
        5.95,
        11.4,
        0.7,
        [("如果只靠传统招商，元谷是“多一个园区”；如果导入政策、圈层和宣传，元谷才有可能成为“上海科技时尚产业场景”。", 15, True, WHITE)],
    )
    footer(s)

    # ========== 04 必要性展开：不是多一个人 ==========
    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "02 · 合作必要性", "我们补的是自有招商最难补齐的四块“基础设施”")
    cards = [
        ("政府通道", "政策申报与对接", "把补贴、资质、算力、租金支持等，前置为入驻谈判筹码。", TEAL),
        ("精准客群", "圈层筛选 + 数据", "企业家/机构/协会/校友高信任客群；爬楼数据压缩无效触达。", ACCENT),
        ("活动获客", "沙龙即招商入口", "6 场 × ≥30 客户；把信任关系变成可跟进线索。", GOLD),
        ("品牌立势", "宣传降低解释成本", "挂牌 + 媒体 + 案例，让客户先相信定位，再谈租金。", PRIMARY2),
    ]
    for i, (t1, t2, t3, ac) in enumerate(cards):
        col, rowi = i % 2, i // 2
        l = 0.7 + col * 6.2
        t = 1.45 + rowi * 2.55
        _rect(s, l, t, 5.95, 2.35, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, l, t, 5.95, 0.12, fill=ac)
        _text(s, l + 0.3, t + 0.35, 5.4, 0.4, [(t1, 18, True, PRIMARY)])
        chip(s, l + 0.3, t + 0.9, 3.6, 0.4, t2, fill=ac, size=12)
        _text(s, l + 0.3, t + 1.5, 5.4, 0.65, [(t3, 13.5, False, MUTED)])
    footer(s)

    # ========== 05 政府政策资源 ==========
    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "03 · 政府政策资源端优势", "把“政策申报能力”转化为企业入驻理由，而不是停留在口头背书")
    steps = [
        ("政策识别", "匹配企业类型", "AI、IP、设计、科技服务、总部/研发中心、数字经济等。", TEAL),
        ("申报辅导", "梳理落地条件", "帮助入驻企业梳理主体、材料、税收、研发、人才与租赁条件。", ACCENT),
        ("政府对接", "打通沟通路径", "对接区级产业部门、科创系统、协会平台与平台招商场景。", GOLD),
        ("结果转化", "变成谈判筹码", "把补贴、荣誉、资质、算力、资本合作转为入驻决策理由。", GREEN),
    ]
    for i, (a, b, c, ac) in enumerate(steps):
        x = 0.7 + i * 3.15
        _rect(s, x, 1.45, 3.0, 3.55, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, x, 1.45, 3.0, 0.7, fill=ac)
        _text(s, x, 1.45, 3.0, 0.7, [(a, 15, True, WHITE)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, x + 0.15, 2.4, 2.7, 0.45, [(b, 14, True, PRIMARY)], align=PP_ALIGN.CENTER)
        _text(s, x + 0.2, 3.05, 2.6, 1.6, [(c, 12.5, False, MUTED)], align=PP_ALIGN.CENTER)
        if i < 3:
            _text(s, x + 2.85, 2.9, 0.35, 0.4, [("→", 18, True, ACCENT)], align=PP_ALIGN.CENTER)
    _rect(s, 0.7, 5.25, 11.95, 1.6, fill=CARD, line=LINE, line_w=0.75)
    _text(s, 0.95, 5.4, 11.5, 0.35, [("对陈总可讲的正式口径", 13, True, PRIMARY)])
    bullets(
        s,
        0.95,
        5.85,
        11.5,
        0.9,
        [
            "我们不是承诺“包拿补贴”，而是把政策口径、申报路径、政府沟通和企业落地诉求前置，帮助森马把政策资源变成招商谈判筹码。",
            "慎用措辞：政策支持以实际申报条件、主管部门审核与项目合规为准。",
        ],
        fs=12.5,
        gap=6,
    )
    footer(s)

    # ========== 06 政策清单 ==========
    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "03 · 政府政策资源端优势", "可围绕的政策与导入方向（示意）")
    table(
        s,
        0.7,
        1.4,
        11.95,
        ["方向", "企业侧价值", "对元谷招商的作用"],
        [
            ["高新技术企业 / 专精特新", "资质背书与专项支持", "吸引有成长诉求的科技企业"],
            ["创新券 / 研发与人才支持", "直接降本、增强留沪意愿", "缩短看房到签约决策周期"],
            ["算力 / 产业基金返投", "资本与算力驱动落地", "服务 AI 主轴客户“带着资源来”"],
            ["文创 / 特色小镇 / 租金相关支持", "属地扶持与落地配套", "强化大零号湾政策叙事"],
            ["总部 / 研发中心落户相关", "空间与政策组合方案", "提高大客户谈判筹码"],
        ],
        col_w=[3.5, 3.9, 4.55],
        row_h=0.7,
        accent_first=True,
    )
    chip(
        s,
        0.7,
        5.5,
        11.95,
        0.7,
        "配合 IP + AI 双轨：AI 侧吃透政策与资本，IP 侧兼顾文化与时尚叙事，避免纯二次元难拿政策的困局。",
        fill=PRIMARY,
        size=13,
    )
    _text(
        s,
        0.7,
        6.4,
        12,
        0.4,
        [("身份前提：上海科协体系相关团队 · 上海市科技企业联合会 · 复旦大学住房政策研究中心（已向森马董事长汇报说明）。", 12, True, ACCENT)],
    )
    footer(s)

    # ========== 07 110:1 漏斗 ==========
    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "04 · 110 : 1 成功客群漏斗", "传统招商是大海捞针；我们的价值是压缩分母、提高每一层转化率")
    layers = [
        ("110", "原始触达", "广撒网 / 中介名单 / 陌生拜访", MUTED),
        ("30", "有效到场/沟通", "愿意见面、愿意听方案", PRIMARY2),
        ("12", "需求匹配", "面积、租金、产业方向匹配", TEAL),
        ("3", "意向谈判", "比政策、比重启成本、比配套", GOLD),
        ("1", "签约落地", "真正落定的成功客群", ACCENT),
    ]
    # funnel bars centered
    widths = [11.2, 9.4, 7.6, 5.8, 4.0]
    y = 1.4
    for (num, name, desc, ac), ww in zip(layers, widths):
        x = 0.7 + (11.95 - ww) / 2
        _rect(s, x, y, ww, 0.78, fill=ac if num in ("1", "3") else CARD, line=LINE if num not in ("1", "3") else None, line_w=0.75)
        if num in ("1", "3"):
            fc, sc = WHITE, RGBColor(0xFF, 0xF0, 0xE0)
        else:
            fc, sc = PRIMARY, MUTED
            _rect(s, x, y, 0.1, 0.78, fill=ac)
        _text(s, x + 0.25, y + 0.08, 1.4, 0.6, [(num, 22, True, fc if num in ("1", "3") else ac)], anchor=MSO_ANCHOR.MIDDLE)
        _text(s, x + 1.7, y + 0.08, 4.5, 0.6, [(name, 14, True, fc)], anchor=MSO_ANCHOR.MIDDLE)
        _text(s, x + 6.0 if ww > 6 else x + 1.7, y + (0.08 if ww > 6 else 0.42), ww - (6.3 if ww > 6 else 0.5), 0.35, [(desc, 11, False, sc)])
        y += 0.88

    three = [
        ("传统招商的痛点", "数量看起来多，但多数无迁址周期、租金能力或政策适配度。", RED),
        ("我们压缩分母", "用政策适配、沙龙筛选、爬楼数据和圈层背书提前过滤。", TEAL),
        ("森马获得结果", "形成可约谈、可复访、可签约的真实客户池。", GREEN),
    ]
    for i, (a, b, ac) in enumerate(three):
        x = 0.7 + i * 4.15
        _rect(s, x, 5.9, 3.95, 1.05, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, x, 5.9, 0.1, 1.05, fill=ac)
        _text(s, x + 0.25, 5.98, 3.5, 0.3, [(a, 12, True, PRIMARY)])
        _text(s, x + 0.25, 6.35, 3.5, 0.5, [(b, 11, False, MUTED)])
    footer(s)

    # ========== 08 漏斗核心表达 ==========
    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "04 · 110 : 1 成功客群漏斗", "核心表达：110:1 不是悲观数字，而是招商规律")
    kpi(s, 0.7, 1.45, 3.85, 1.7, "行业规律", "110 → 1", "触达一百多家，往往才签一家", accent=RED)
    kpi(s, 4.75, 1.45, 3.85, 1.7, "合作要做的事", "110 → 30", "把原始触达压成高质量有效沟通", accent=ACCENT)
    kpi(s, 8.8, 1.45, 3.85, 1.7, "最终目标", "30 → 签约", "提高复访率与签约率", accent=GREEN)
    _rect(s, 0.7, 3.45, 11.95, 3.4, fill=CARD, line=LINE, line_w=0.75)
    _text(s, 1.0, 3.7, 11.4, 0.4, [("合作必要性的招商语言", 15, True, PRIMARY)])
    bullets(
        s,
        1.0,
        4.3,
        11.4,
        2.3,
        [
            "110:1 说明：普通招商一定会消耗大量时间与无效沟通成本。",
            "我们的价值：用资源体系把 110 变成更高质量的 30，再把 30 推向可签约。",
            "对森马：不是“多招几个人”，而是少走弯路——让每一次到访都更接近成交。",
            "验证方式：先用首场沙龙与首批客户复访，检验真实转化，再放大合作。",
        ],
        fs=14,
        gap=10,
    )
    footer(s)

    # ========== 09 沙龙客群 ==========
    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "05 · 沙龙带来的优质客群", "沙龙不是活动成本，而是把信任关系变成招商线索的入口")
    groups = [
        ("企业家 / 创始人", "AI、IP、文创、科技服务、设计研发类企业，具备迁址与扩张可能。", ACCENT),
        ("投资机构 / 银行", "基金、银行、产业资本，可反向带项目、带客户、带融资场景。", TEAL),
        ("政府 / 协会平台", "产业部门、科协/科企联、商协会，形成政策与信用背书。", GOLD),
        ("高校 / 校友圈层", "复旦、北大等校友与专家资源，增强专业可信度与传播性。", PRIMARY2),
    ]
    for i, (a, b, ac) in enumerate(groups):
        col, rowi = i % 2, i // 2
        l = 0.7 + col * 6.2
        t = 1.4 + rowi * 2.0
        _rect(s, l, t, 5.95, 1.8, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, l, t, 0.12, 1.8, fill=ac)
        _text(s, l + 0.35, t + 0.3, 5.3, 0.4, [(a, 16, True, PRIMARY)])
        _text(s, l + 0.35, t + 0.85, 5.3, 0.7, [(b, 13, False, MUTED)])
    footer(s)

    # ========== 10 沙龙量化 ==========
    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "05 · 沙龙带来的优质客群", "可计量的客群管道，而不是“热闹一场”")
    kpi(s, 0.7, 1.4, 3.85, 1.55, "年度沙龙", "6 场", "主题招商，不是泛活动", accent=ACCENT)
    kpi(s, 4.75, 1.4, 3.85, 1.55, "年度精准触达", "180+ 家", "6 场 × ≥30 客户", accent=TEAL)
    kpi(s, 8.8, 1.4, 3.85, 1.55, "质量对比", "高于泛推", "客群质量高于普通中介名单", accent=GREEN)
    table(
        s,
        0.7,
        3.2,
        11.95,
        ["场次主题（示意）", "联办资源", "客群画像"],
        [
            ["AI + 潮玩 / 科技时尚", "产业平台 + 科技资源", "AI / 内容科技中大型"],
            ["潮玩出海", "国际会客厅 + 影响力平台", "出海 IP / 外贸型"],
            ["投融资路演", "产业基金 + 银行", "有融资诉求的成长企业"],
            ["设计与创意", "高校 + 科企联", "设计总部 / 创意机构"],
            ["内容 IP · Z 世代", "行业协会 + 产业平台", "内容、直播、IP 运营"],
            ["政策补贴 · 小镇", "区级科创系统 + 学术平台", "政策敏感型科技企业"],
        ],
        col_w=[3.6, 4.0, 4.35],
        row_h=0.48,
        accent_first=True,
    )
    footer(s)

    # ========== 11 沙龙转化路径 ==========
    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "05 · 沙龙带来的优质客群", "从沙龙到签约：四步转化路径")
    path = [
        ("场前邀约", "定向名单", "科企联资源池 + 精准客户清单 + 峰会/榜单 VIP", ACCENT),
        ("场中触达", "≥30 家/场", "主题匹配，天然过滤无效流量", TEAL),
        ("场后复访", "意向池", "到访元谷、政策匹配、方案比选", GOLD),
        ("签约转化", "可验证结果", "进入正式租赁谈判与落定", GREEN),
    ]
    for i, (a, b, c, ac) in enumerate(path):
        x = 0.7 + i * 3.15
        _rect(s, x, 1.5, 3.0, 3.4, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, x, 1.5, 3.0, 0.55, fill=ac)
        _text(s, x, 1.5, 3.0, 0.55, [(a, 14, True, WHITE)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, x + 0.15, 2.35, 2.7, 0.55, [(b, 18, True, PRIMARY)], align=PP_ALIGN.CENTER)
        _text(s, x + 0.2, 3.15, 2.6, 1.4, [(c, 13, False, MUTED)], align=PP_ALIGN.CENTER)
        if i < 3:
            _text(s, x + 2.85, 2.9, 0.35, 0.4, [("→", 18, True, ACCENT)], align=PP_ALIGN.CENTER)
    bullets(
        s,
        0.7,
        5.2,
        12,
        1.6,
        [
            "峰会 / 榜单 / 闭门会可形成二次复访，把“见过一面”变成“反复沟通”。",
            "单场媒体曝光可同步发生，客群获客与品牌声量并行。",
            "对资管侧：成果可追踪——名单、跟进、到访、签约，便于检验合作价值。",
        ],
        fs=13.5,
        gap=8,
    )
    footer(s)

    # ========== 12 宣传作用 ==========
    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "06 · 我们对项目的宣传作用", "把元谷从“可出租空间”变成“可被报道、可被政府看见、可被企业引用的产业事件”")
    phases = [
        ("会前", "主题包装", "AI+IP / 科技时尚 / 政策申报服务包；媒体稿、案例文章、榜单与挂牌预热。", ACCENT),
        ("会中", "现场转化", "专属展位、闭门洽谈、企业需求登记；把声量当场变成线索。", TEAL),
        ("会后", "内容传播", "媒体稿、案例文章、榜单/挂牌、客户复访；持续降低后续解释成本。", GOLD),
    ]
    for i, (a, b, c, ac) in enumerate(phases):
        x = 0.7 + i * 4.15
        _rect(s, x, 1.45, 3.95, 3.5, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, x, 1.45, 3.95, 0.95, fill=ac)
        _text(s, x + 0.2, 1.55, 3.55, 0.35, [(a, 12, False, RGBColor(0xFF, 0xF0, 0xE0))], align=PP_ALIGN.CENTER)
        _text(s, x + 0.2, 1.9, 3.55, 0.4, [(b, 18, True, WHITE)], align=PP_ALIGN.CENTER)
        _text(s, x + 0.25, 2.7, 3.45, 2.0, [(c, 13.5, False, TEXT)])
    _rect(s, 0.7, 5.2, 11.95, 1.65, fill=PRIMARY)
    _text(s, 1.0, 5.4, 11.4, 0.35, [("对陈总的价值", 13, True, RGBColor(0xE8, 0xC8, 0x7A))])
    _text(
        s,
        1.0,
        5.9,
        11.4,
        0.75,
        [("宣传不是“面子工程”，而是降低招商解释成本。客户看到政府、协会、专家、媒体共同出现，会更容易相信项目定位与未来产业氛围。", 14.5, True, WHITE)],
    )
    footer(s)

    # ========== 13 宣传闭环 ==========
    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "06 · 我们对项目的宣传作用", "宣传闭环：活动造势 × 媒体放大 × 招商转化")
    ring = [
        ("① 挂牌立势", "科企联 / 学术 / 产业平台牌照", "公信力", PRIMARY),
        ("② 沙龙造势", "6 场产业主题活动", "客群", ACCENT),
        ("③ 媒体放大", "主流媒体 + 社交投放", "声量", TEAL),
        ("④ 招商转化", "名单跟进 → 到访 → 签约", "结果", GREEN),
    ]
    for i, (a, b, c, ac) in enumerate(ring):
        x = 0.7 + i * 3.15
        _rect(s, x, 1.45, 3.0, 2.55, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, x, 1.45, 3.0, 0.55, fill=ac)
        _text(s, x, 1.45, 3.0, 0.55, [(a, 13, True, WHITE)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, x + 0.15, 2.25, 2.7, 0.9, [(b, 13, False, TEXT)], align=PP_ALIGN.CENTER)
        chip(s, x + 0.55, 3.4, 1.9, 0.35, c, fill=PRIMARY, size=11)
    table(
        s,
        0.7,
        4.3,
        11.95,
        ["可选媒体模块（示意）", "量级", "说明"],
        [
            ["原创内容（品牌/营销稿）", "3 篇", "媒体总曝光预期 ≥150 万（可按服务包执行）"],
            ["主流媒体宣发", "约 10 篇", "中央级 / 全国 / 上海媒体档位据实"],
            ["社交信息流投放", "按量", "抖音 / 头条等，上海区域精准"],
        ],
        col_w=[4.5, 2.2, 5.25],
        row_h=0.55,
        accent_first=True,
    )
    _text(
        s,
        0.7,
        6.4,
        12,
        0.4,
        [("说明：媒体为可选模块，可与沙龙、挂牌独立决策；核心是用宣传服务招商结果，而不是为热闹买单。", 12.5, True, ACCENT)],
    )
    footer(s)

    # ========== 14 三步走 ==========
    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "07 · 合作落地建议", "三步走：先用低成本启动资源闭环，再用签约结果证明价值")
    steps = [
        ("第 1 步", "确认合作口径", "本周", "确定项目定位、授权边界、宣传口径、政策服务边界。", ACCENT),
        ("第 2 步", "启动首场沙龙", "2–3 周", "围绕 AI+IP / 科技时尚 / 政策申报，邀请首批 30+ 有效客户。", TEAL),
        ("第 3 步", "客户复访与签约", "30–90 天", "形成意向客户清单、政策匹配表、复访计划、招商谈判排期。", GREEN),
    ]
    for i, (no, title, when, body, ac) in enumerate(steps):
        x = 0.7 + i * 4.15
        _rect(s, x, 1.45, 3.95, 4.0, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, x, 1.45, 3.95, 1.25, fill=ac)
        _text(s, x + 0.2, 1.55, 3.55, 0.35, [(no, 12, False, RGBColor(0xFF, 0xF0, 0xE0))], align=PP_ALIGN.CENTER)
        _text(s, x + 0.2, 1.95, 3.55, 0.55, [(title, 18, True, WHITE)], align=PP_ALIGN.CENTER)
        chip(s, x + 0.85, 2.95, 2.2, 0.4, when, fill=PRIMARY, size=12)
        _text(s, x + 0.25, 3.6, 3.45, 1.5, [(body, 13.5, False, TEXT)])
    _text(
        s,
        0.7,
        5.7,
        12,
        1.1,
        [
            ("建议给陈总的底线", 13, True, PRIMARY),
            ("森马不需要先相信所有结果，但要先启动一个可验证的合作试点；用首场沙龙 + 首批客户复访来检验真实转化。", 14, True, ACCENT),
        ],
    )
    footer(s)

    # ========== 15 价值总述（外发口径） ==========
    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "08 · 合作价值总述", "我们不是来卖活动，也不是简单做中介")
    _rect(s, 0.7, 1.4, 11.95, 2.55, fill=PRIMARY)
    _text(s, 1.0, 1.6, 11.4, 0.4, [("正式价值主张", 13, True, RGBColor(0xE8, 0xC8, 0x7A))])
    _text(
        s,
        1.0,
        2.15,
        11.4,
        1.5,
        [
            ("帮助森马把“空间出租”升级为“政策 + 圈层 + 宣传 + 招商”的闭环。", 18, True, WHITE),
            ("用外部资源降低获客成本、提高客户信任、加快签约转化。", 15, False, RGBColor(0xCF, 0xD8, 0xE8)),
        ],
    )
    three = [
        ("① 政策与申报路径", "把企业入驻理由做实，合规前置、路径清晰。", TEAL),
        ("② 高质量圈层客群", "企业家、机构、协会、校友，把泛流量变成精准客户。", ACCENT),
        ("③ 产业场景包装", "把元谷做成上海科技时尚产业场景，降低认知成本。", GOLD),
    ]
    for i, (a, b, ac) in enumerate(three):
        x = 0.7 + i * 4.15
        _rect(s, x, 4.2, 3.95, 2.05, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, x, 4.2, 3.95, 0.12, fill=ac)
        _text(s, x + 0.25, 4.5, 3.45, 0.45, [(a, 15, True, PRIMARY)])
        _text(s, x + 0.25, 5.15, 3.45, 0.85, [(b, 13, False, MUTED)])
    footer(s)

    # ========== 16 结论 ==========
    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "结论", "请陈总今天带走的四个判断")
    conclusions = [
        ("必要性", "不是替代资管，而是补齐政策、圈层、宣传三条腿；否则 110:1 的成本由项目自己扛。", ACCENT),
        ("政策资源", "把政策申报能力变成入驻理由；合规前置，不做“包拿”承诺。", TEAL),
        ("沙龙客群", "6 场 × ≥30 = 180+ 精准触达；企业家/机构/协会/校友高信任客群。", GOLD),
        ("宣传作用", "把元谷从出租空间升级为可报道的产业事件，降低招商解释成本。", PRIMARY2),
    ]
    for i, (t1, t2, ac) in enumerate(conclusions):
        y = 1.4 + i * 1.2
        _rect(s, 0.7, y, 11.95, 1.05, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, 0.7, y, 2.3, 1.05, fill=ac)
        _text(s, 0.7, y, 2.3, 1.05, [(t1, 16, True, WHITE)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, 3.25, y + 0.25, 9.1, 0.6, [(t2, 14, False, TEXT)], anchor=MSO_ANCHOR.MIDDLE)
    footer(s)

    # ========== 17 结束页 ==========
    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, 13.333, 7.5, fill=PRIMARY)
    _rect(s, 0, 0, 0.28, 7.5, fill=ACCENT)
    _text(s, 1.1, 2.0, 11, 0.4, [("下一步", 14, False, RGBColor(0xE8, 0xC8, 0x7A))])
    _text(s, 1.1, 2.55, 11, 1.0, [("先启动可验证试点", 30, True, WHITE), ("再用结果决定放大节奏", 30, True, WHITE)])
    bullets(
        s,
        1.1,
        4.5,
        11,
        1.8,
        [
            "本周：确认合作口径与授权边界；",
            "2–3 周：启动首场沙龙（30+ 有效客户）；",
            "30–90 天：客户复访、政策匹配、签约推进。",
        ],
        fs=15,
        gap=10,
        color=RGBColor(0xE5, 0xEB, 0xF5),
        marker="▸",
        mc=RGBColor(0xE8, 0xC8, 0x7A),
    )
    _text(
        s,
        1.1,
        6.6,
        11,
        0.4,
        [("胡教授团队｜复旦大学住房政策研究中心 / 上海市科技企业联合会相关资源", 12, False, RGBColor(0x9F, 0xB2, 0xD0))],
    )

    prs.save(OUT)
    print(f"Wrote {OUT} | content pages counted in footer: {_page['n']}")


if __name__ == "__main__":
    build()
