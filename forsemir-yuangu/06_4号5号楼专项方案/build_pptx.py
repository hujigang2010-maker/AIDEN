"""Build the 4#+5# 楼 2 万方专项 PPT (.pptx).

九段式: 执行摘要 → 项目实测 → SWOT → Phase 1-5 → 投决建议.
关键资产: AI 潮玩产业基地 + 6 场产业沙龙 + 5 项挂牌 + 5 月 22 日 AI 峰会 + 仲量联行爬楼大数据 + 追觅基金.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

OUT = Path(__file__).with_name("4#5#楼产业招商运营方案.pptx")

PRIMARY = RGBColor(0x14, 0x2C, 0x5E)
ACCENT = RGBColor(0xF2, 0x7E, 0x2D)
GOLD = RGBColor(0xC8, 0x99, 0x3D)
LIGHT = RGBColor(0xF5, 0xF7, 0xFA)
TEXT = RGBColor(0x1F, 0x2A, 0x44)
MUTED = RGBColor(0x55, 0x60, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.95))
    bar.fill.solid(); bar.fill.fore_color.rgb = PRIMARY; bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.4); tf.margin_top = Inches(0.15)
    tf.text = title
    p = tf.paragraphs[0]; p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = WHITE
    if subtitle:
        sub = tf.add_paragraph(); sub.text = subtitle
        sub.font.size = Pt(12); sub.font.color.rgb = RGBColor(0xCF, 0xD8, 0xE3)


def add_body(slide, bullets, left=0.6, top=1.25, width=12.1, height=5.7, font_size=16):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame; tf.word_wrap = True
    for i, line in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        depth = 0; text = line
        if text.startswith("    "):
            depth = 2; text = text.lstrip()
        elif text.startswith("  "):
            depth = 1; text = text.lstrip()
        para.text = text
        para.level = depth
        para.font.size = Pt(font_size if depth == 0 else font_size - 2)
        para.font.color.rgb = TEXT if depth == 0 else MUTED
        para.font.bold = depth == 0 and (text.startswith("【") or text.startswith("★"))
        para.space_after = Pt(5)


def add_footer(slide, text="胡教授团队 × 森马集团 · 元谷 4#+5# 楼 2 万方专项招商运营方案"):
    box = slide.shapes.add_textbox(Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.4))
    tf = box.text_frame; tf.text = text
    p = tf.paragraphs[0]; p.font.size = Pt(10); p.font.color.rgb = MUTED


def add_table(slide, top_inch, headers, rows, left=0.6, width=12.1, header_color=None):
    cols = len(headers); n = len(rows) + 1
    header_color = header_color or PRIMARY
    table_shape = slide.shapes.add_table(n, cols, Inches(left), Inches(top_inch), Inches(width), Inches(0.5 + 0.42 * len(rows)))
    table = table_shape.table
    for j, h in enumerate(headers):
        cell = table.cell(0, j); cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = WHITE; r.font.bold = True; r.font.size = Pt(12)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j); cell.text = val
            cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT if i % 2 else WHITE
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.color.rgb = TEXT; r.font.size = Pt(10)


def add_kpi_card(slide, left_in, top_in, w_in, h_in, headline, value, sub):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_in), Inches(top_in), Inches(w_in), Inches(h_in))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT; box.line.color.rgb = PRIMARY
    tf = box.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.15)
    tf.text = headline
    p = tf.paragraphs[0]; p.font.size = Pt(11); p.font.color.rgb = MUTED; p.font.bold = False
    p2 = tf.add_paragraph(); p2.text = value
    p2.font.size = Pt(20); p2.font.color.rgb = PRIMARY; p2.font.bold = True
    p3 = tf.add_paragraph(); p3.text = sub
    p3.font.size = Pt(10); p3.font.color.rgb = MUTED


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ============== P1 Cover ==============
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = PRIMARY; bg.line.fill.background()
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(3.0), Inches(0.3), Inches(2.4))
    accent.fill.solid(); accent.fill.fore_color.rgb = ACCENT; accent.line.fill.background()
    tf = s.shapes.add_textbox(Inches(1.1), Inches(2.4), Inches(11.5), Inches(3.5)).text_frame
    tf.text = "元谷项目"
    tf.paragraphs[0].font.size = Pt(56); tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = WHITE
    p = tf.add_paragraph(); p.text = "4#+5# 楼约 2 万㎡ 产业研发办公"
    p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = WHITE
    p2 = tf.add_paragraph(); p2.text = "招 商 运 营 专 项 方 案"
    p2.font.size = Pt(28); p2.font.color.rgb = ACCENT
    p3 = tf.add_paragraph(); p3.text = ""
    p4 = tf.add_paragraph(); p4.text = "胡教授团队  ×  森马集团     ·     202X 年提案 v1.0"
    p4.font.size = Pt(16); p4.font.color.rgb = RGBColor(0xCF, 0xD8, 0xE3)

    # ============== P2 执行摘要 ==============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P2 · 执行摘要", "Executive Summary — 30 秒读懂本方案")
    add_body(s, [
        "★ 业务包：元谷 4# / 5# 楼 5F+ 共约 2 万㎡产业研发办公整体招商运营 (24 个月期)",
        "★ 核心抓手：『AI 潮玩产业基地』 + 6 场产业沙龙 + 5 项挂牌 + 5/22 AI 商业化峰会借势",
        "★ 团队配置：2 人驻场 + CSO 顾问 (轻骑兵) → 月度成本可控",
        "★ 商业条款 (向森马提议)：",
        "  ① 基础月费 (Retainer) 12 万元/月 (2 人推荐配)",
        "  ② 招商佣金 1.5 / 1.75 / 2.0 个月年租金 (按面积阶梯)",
        "  ③ 挂牌奖励 30 万元/项 × 5 项 = 150 万元一次性",
        "  ④ 沙龙执行 5 万元/场 × 6 场 + 净利 30/70 分润",
        "★ 24 个月给森马的回报：满租后年租金 ≈ 1,606 万元 + 物业 ≈ 240 万 + 资产增值",
        "★ 投入产出 (甲方视角)：付出 818-965 万 → 获得 3,200 万+ 直接年化收入 (ROI ≈ 1:3.9)",
    ], font_size=15)
    add_footer(s)

    # ============== P3 项目实测 (一) 楼宇结构 ==============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P3 · 项目实测 (一)：4# / 5# 楼业态结构", "本方案聚焦 5F+ 产业集群约 2 万㎡")
    add_table(s, 1.2,
        ["楼栋", "1–3F", "4F", "5F+", "本方案招商范围"],
        [
            ["4# 楼", "潮玩艺术中心", "直播中心", "潮玩产业集群 ≈ 1 万㎡", "★ 国际 IP 创意层"],
            ["5# 楼", "动漫书店 / 休闲娱乐", "动漫书店延伸", "潮玩产业集群 ≈ 1 万㎡ (含 5F 潮玩产业展厅)", "★ 潮玩产业总部层"],
        ])
    add_body(s, [
        "4# 楼 5F+ 已具备的同栋配套:4F 直播中心 / 5F AI 共享设计中心 / 5F AI 共享打样 & DIY 中心",
        "5# 楼 5F+ 已具备的同栋配套:5F 潮玩产业展厅 / 1–4F 动漫书店及休闲娱乐 (天然客流和场景)",
        "→ 4#/5# 楼 5F+ 是元谷整盘的产业心脏, 与 1F-4F 商业业态形成『前店后办公』复合形态",
    ], top=3.7, height=3.0, font_size=14)
    add_footer(s)

    # ============== P4 项目实测 (二) 区位与租金 ==============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P4 · 项目实测 (二)：区位、客流与租金", "与大零号湾主流园区对标")
    add_kpi_card(s, 0.6, 1.3, 3.0, 1.2, "TOD 加持", "15 号线元江路站", "单日客流 5–7 万人次")
    add_kpi_card(s, 3.7, 1.3, 3.0, 1.2, "15 分钟车行覆盖", "24 万 居住", "+ 12 万 产业办公")
    add_kpi_card(s, 6.8, 1.3, 3.0, 1.2, "市场对标日租金", "1.8–2.5 元/㎡/天", "(大零号湾 / 紫竹高新区)")
    add_kpi_card(s, 9.9, 1.3, 3.0, 1.2, "建议成交日租金", "2.0–2.5 元", "中位 2.2 (含 5A 加持)")
    add_table(s, 2.8,
        ["对标园区", "日租金 (元/㎡/天)", "类别"],
        [
            ["零号湾全球创新创业集聚区 / 大零号湾科创转化中心", "2.0–2.5", "标杆"],
            ["华谊万创新所 / 上海人工智能产业园", "2.2", "主流"],
            ["紫竹信息数码港 (5A 甲级) / 紫竹数字创意港", "2.1–3.0", "主流"],
            ["龙湖蓝海引擎 / 金领谷科技产业园", "1.5–2.8", "可对标"],
            ["云境 443 未来产业社区 / 夏日汇国际中心 (高端)", "2.3–4.5", "上限"],
        ])
    add_body(s, [
        "满租推算 (甲方视角): 2 万㎡ × 365 天 × 2.2 元 ≈ 1,606 万元/年租金 (永续年金)",
    ], top=5.6, height=1.0, font_size=14)
    add_footer(s)

    # ============== P5 项目实测 (三) 与全盘的关系 ==============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P5 · 项目实测 (三)：与 22 万㎡ 全盘的关系")
    add_body(s, [
        "【本方案聚焦】4#+5# 楼 5F+ 共 2 万㎡产业研发办公 → 元谷整盘的『产业心脏』",
        "【整盘 22 万㎡背景，仅作概览】",
        "  · 1# 楼: 森马集团总部办公 + 1-4F 零售",
        "  · 2# 楼: 二次元主题 Livehouse / 秀场",
        "  · 3# 楼: 1-4F 休闲运动 + 萌宠 / 5F+ 酒店",
        "  · 4# 楼: 1-3F 潮玩艺术中心 / 4F 直播中心 / ★ 5F+ 潮玩产业集群 (本方案)",
        "  · 5# 楼: 1-4F 动漫书店 + 休闲娱乐 / ★ 5F+ 潮玩产业集群 (本方案)",
        "  · 6# 楼: 1-5F 餐饮 / 服务配套 / 商务宴请",
        "【商业 5.2 万㎡部分】属于二期拓展空间，本阶段不在专项报告范围内 (待 4#/5# 成势后承接)",
        "【为什么选 4#/5# 5F+】产业属性最纯粹、配套最齐全、单价容易守住、可以快速锁定 AI 潮玩产业基地双牌照",
    ], font_size=15)
    add_footer(s)

    # ============== P6 SWOT ==============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P6 · SWOT 矩阵")
    add_table(s, 1.3,
        ["", "内部 (Internal)", "外部 (External)"],
        [
            ["有利",
             "S 优势 · 森马底盘 + 双牌照已确认 + 4#/5# 楼产业配套齐全 (选品/直播/AI 设计/展厅/艺术中心)",
             "O 机会 · AI+潮玩双赛道政策红利 / 5/22 峰会高净值人脉 / 复旦+北大学术资源可平移"],
            ["不利",
             "W 劣势 · 新建项目品牌势能尚弱 / TOD 工程未完工 / 潮玩生态需培育",
             "T 威胁 · 周边新园区 (云境 443 / 夏日汇) 竞争 / 政策调整 / 招商节奏不及预期"],
        ])
    add_body(s, [
        "【应对策略】",
        "  · 用『基础月费保底 + 阶梯佣金』对冲 W (招商节奏风险)",
        "  · 用『5/22 峰会 + 6 场沙龙 + 5 项挂牌』转化 O (政策红利) 为实际客户",
        "  · 用 AI 潮玩产业基地『差异化定位』对冲 T (周边竞争)",
    ], top=4.4, height=2.4, font_size=14)
    add_footer(s)

    # ============== P7 Phase 1 总定位 ==============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P7 · Phase 1 · 策略与定位 (一)", "AI 潮玩产业基地 — 长三角首个 AI+潮玩 双牌照策源高地")
    add_body(s, [
        "★ 总定位:『AI 潮玩产业基地 (China AI + IP-Toy Industrial Base)』",
        "  — 长三角首个 AI + 潮玩 双牌照产业策源高地",
        "  — 以中国动漫集团『AI 潮玩产业基地』+ 中国百货商业协会『潮玩次元商业专委会』为核心牌照",
        "  — 与上海交大设计学院 / 闵行科协 / 森马集团 联合建设",
        "【三大战略目标】",
        "  ① 国际化产业枢纽 (4# 楼 5F+, 北欧 / 日韩 / 东南亚 IP 首站)",
        "  ② 品牌运营高地 (5# 楼 5F+, 头部企业 + 中型潮玩运营企业)",
        "  ③ IP 创制中心 (4# 楼 5F AI 共享设计 + 5# 楼 5F 潮玩产业展厅)",
        "【与森马原产业规划完全对齐】PDF 中『科技潮玩产业策源高地』+『五档配比 10/10/20/20/40』 → 本方案直接承接执行",
    ], font_size=15)
    add_footer(s)

    # ============== P8 Phase 1 楼宇分工 ==============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P8 · Phase 1 · 策略与定位 (二)：4# / 5# 楼分工")
    add_table(s, 1.2,
        ["定位维度", "4# 楼 5F+ (≈ 1 万㎡)", "5# 楼 5F+ (≈ 1 万㎡)"],
        [
            ["昵称",                     "国际 IP 创意层",                "潮玩产业总部层"],
            ["核心客户",                  "国际 IP / AI 设计 / 国际艺术家",  "头部央企 / 行业协会 / 中型潮玩运营"],
            ["同层配套联动",               "4F 直播中心 / 5F AI 共享设计中心", "5F 潮玩产业展厅 / 1-4F 动漫书店"],
            ["重点 IP 平台",               "北欧创新国际会客厅 元谷站",       "AI 潮玩产业基地 主基地"],
            ["典型户型",                  "200–800㎡ 小型 + 1000–2000㎡ 中型", "300–1,200㎡ 中型 + 2,000㎡ 头部"],
            ["首年优先",                  "★ T+1 年优先满租 (1 万㎡)",       "T+2 年承接外溢 + 满租 (1 万㎡)"],
        ])
    add_footer(s)

    # ============== P9 Phase 1 客户画像 ==============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P9 · Phase 1 · 策略与定位 (三)：五档客户画像 (与 PDF 一致)")
    add_table(s, 1.2,
        ["占比", "类型", "数量 / 户型", "落位", "招商打法"],
        [
            ["10%", "头部央企 / 行业协会 (导向)",   "3 家 / 2,000㎡",      "5# 楼 5F+",     "L1 牌照锚定"],
            ["10%", "共享配套服务体系 (吸附点)",     "3 家 / 2,000㎡",      "4# 楼 5F+ + 5# 楼 5F+", "合资公司直营"],
            ["20%", "中型潮玩运营企业 (基础)",       "4-6 家 / 5,000㎡",   "5# 楼 5F+",     "L2 资本招商 (返投基金)"],
            ["20%", "中小型潮业服务机构 (血肉)",      "15 家 / 200-500㎡", "4# 楼 5F+",     "L4 活动带流 (6 场沙龙)"],
            ["40%", "小型潮玩运营企业 (骨架)",        "30 家 / 200-500㎡", "4# 楼 5F+ 主力", "L3 爬楼大数据"],
        ])
    add_body(s, [
        "★ 总计入驻企业 55+ 家、总人口约 1,000-1,500 人 → 元谷产业人口主战场",
    ], top=4.5, height=0.8, font_size=14)
    add_footer(s)

    # ============== P10 Phase 2 招商漏斗 ==============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P10 · Phase 2 · 招商执行 (一)：四级招商漏斗")
    add_table(s, 1.2,
        ["层级", "招商动作", "核心资源", "目标产能", "周期"],
        [
            ["L1 牌照锚定",   "签头部央企 / 行业协会作为旗舰",        "AI 潮玩产业基地 + 潮玩次元商业专委会 双牌照", "3 家 × 2,000㎡ = 6,000㎡", "T+3 月"],
            ["L2 资本招商",   "追觅基金返投 + AI 腾讯生态导流",       "追觅基金 + 5/22 峰会 LP 资源",                "4-6 家 × 1,000㎡ = 5,000㎡", "T+6 月"],
            ["L3 大数据爬楼", "仲量联行爬楼大数据精准画像 + 上门拜访", "仲量联行爬楼数据 (¥2.6 万已购入)",            "30 家 × 200-500㎡ = 6,000-8,000㎡", "T+9 月"],
            ["L4 活动带流",   "6 场产业沙龙 + 福布斯榜单 + 北欧外事",   "5/22 峰会 + 6 场沙龙",                          "15 家 × 200-500㎡ = 3,000-4,000㎡", "T+12 月"],
        ])
    add_body(s, [
        "★ 漏斗共能产生 20,000–23,000㎡ 招商管道 → 实际签约可达 18,000-20,000㎡ (即满租)",
    ], top=4.6, height=0.8, font_size=14)
    add_footer(s)

    # ============== P11 Phase 2 5.22 峰会借势 ==============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P11 · Phase 2 · 招商执行 (二)：5/22 AI 峰会借势路径", "2026 人工智能商业化落地与硬核投资破局峰会")
    add_body(s, [
        "【峰会概览】2026 年 5 月 22 日 · 上海·北外滩·一滴水 · 13:00–20:30",
        "  主办: 北京大学经济学院上海校友会 + 复旦大学住房政策研究中心",
        "  联办: 上海杨浦区科技企业联合会 / 上海市临港区科技企业联合会 / 北大金融校友联合会 / 中东投资联盟 等",
        "  四大篇章 + VIP 高端闭门晚宴 → 嘉宾覆盖 AI / 量子 / 金融 / 资本 / 文创 / 出海 全领域",
        "【借势路径 — 五件事】",
        "  ① 元谷设峰会专属展位 + 闭门 1V1 招商台 (现场锁定 30-50 位 LP / 潜在客户)",
        "  ② 颁奖嵌入 — 增设『AI 潮玩产业影响力榜』, 由元谷主办、福布斯背书",
        "  ③ 主办方资源平移 — 复旦住房政策中心、北大经济学院上海校友会、中行 / 长江 / 招行 / 金浦 / 铂帝 战略合作",
        "  ④ 峰会嘉宾转化沙龙 #1 主题嘉宾 — 直接转化为元谷 T+1 月首场 AI+潮玩沙龙嘉宾",
        "  ⑤ 峰会闭门晚宴 → 形成『元谷创始合伙人圈』 → 12 个月后归入元谷年度大会",
        "★ 借势效益: 借助峰会高净值人群 (200+ VIP) , 把品牌势能 + 客户线索 提前到 T+0 之前",
    ], font_size=14)
    add_footer(s)

    # ============== P12 Phase 2 客户清单(脱敏) ==============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P12 · Phase 2 · 招商执行 (三)：意向客户管道 (示例 / 脱敏)")
    add_table(s, 1.2,
        ["客户类型", "代表企业 (脱敏)", "意向面积", "对接渠道"],
        [
            ["L1 头部央企",      "中字头数字央企 A / B",           "2,000㎡ × 2",     "中国动漫集团 + 牌照锚定"],
            ["L1 行业协会",      "中国百货商业协会 潮玩次元专委",   "1,500㎡",         "牌照共建"],
            ["L2 中型潮玩 / AI", "AI 潮玩品牌 X / 国漫公司 Y",     "1,000㎡ × 4-6",  "追觅基金 + AI 腾讯"],
            ["L3 小型潮玩",     "盲盒新锐 / 设计师工作室",         "200-500㎡ × 30", "仲量联行爬楼大数据 + 5/22 峰会"],
            ["L4 服务机构",     "财税 / IP 律所 / 出海咨询",      "200-500㎡ × 15", "6 场沙龙拉新"],
        ])
    add_body(s, [
        "★ 实际清单超过 60 家 (覆盖 4#+5# 楼 2 万㎡), 详见 Excel Sheet '招商客户管道'",
    ], top=4.6, height=0.8, font_size=14)
    add_footer(s)

    # ============== P13 Phase 2 招商里程碑 ==============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P13 · Phase 2 · 招商执行 (四)：12+12 个月招商里程碑")
    add_table(s, 1.2,
        ["时间节点", "招商累计 (㎡)", "里程碑事件", "佣金累计 (按 1.75 个月平均, 万元)"],
        [
            ["T+3 月",  "4,000",   "首批牌照锚定 (L1) + 5/22 峰会借势完成", "47"],
            ["T+6 月",  "8,000",   "中型潮玩 4 家入驻 (L2 资本招商)",       "94"],
            ["T+9 月",  "12,000",  "爬楼大数据 (L3) 转化 + 沙龙 4 场完成",   "141"],
            ["T+12 月", "10,000",  "4# 楼 1 万方满租 (L3+L4 双轮)",          "117 (4# 楼全部)"],
            ["T+18 月", "14,000",  "5# 楼 5F+ 启动 + 沙龙转固定 IP",        "164"],
            ["T+24 月", "20,000",  "★ 4#+5# 楼共 2 万方满租 + 沙龙 IP 化",   "234 (累计)"],
        ])
    add_body(s, [
        "★ 24 个月累计佣金 ≈ 234 万元 (基础场景, 平均 1.75 个月年租金)",
    ], top=4.7, height=0.8, font_size=14)
    add_footer(s)

    # ============== P14 Phase 2 招商出彩点 ==============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P14 · Phase 2 · 招商执行 (五)：六大『出彩点』(VS 普通招商)")
    add_body(s, [
        "①『牌照即招商』— AI 潮玩产业基地 + 潮玩次元专委 双牌照前置, 客户『送进来』不是『拉进来』",
        "②『基金即招商』— 追觅基金 1:1.5 返投绑定, 资本驱动签约 (不只是租约, 还有股权)",
        "③『数据即招商』— 仲量联行爬楼大数据 (¥2.6 万已购) → 200 家精准客户清单 + 转化率提升 30%",
        "④『峰会即招商』— 5/22 AI 峰会 200+ VIP, 1 天锁定 30-50 个客户, 普通园区一个季度才够",
        "⑤『沙龙即招商』— 6 场产业沙龙, 每场 30+ 目标客户, 一场签 1-2 家 = 5-12 家直接成果",
        "⑥『学术即招商』— 复旦住房政策研究中心 + 上海交大设计学院 + 北大经济学院 三所高校背书招商",
        "",
        "→ 普通园区只做 ③④, 元谷六维齐发 → 招商速度领先 6-12 个月、单方租金溢价 0.2-0.4 元/㎡/天",
    ], font_size=15)
    add_footer(s)

    # ============== P15 Phase 3 品牌+活动 (一) 6 场沙龙 ==============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P15 · Phase 3 · 品牌与活动 (一)：6 场产业沙龙 (每场 ≥ 30 目标产业客户)")
    add_table(s, 1.2,
        ["序号", "时间", "主题", "联办方", "目标客户数"],
        [
            ["#1", "T+1 月 (借势 5/22 峰会次月)", "AI + 潮玩 跨界融合",         "中国动漫集团 + AI 腾讯",                "30+"],
            ["#2", "T+3 月",                    "潮玩出海 (北欧 / 日韩 / 东南亚)", "北欧创新国际会客厅 + 福布斯",            "30+"],
            ["#3", "T+5 月",                    "投融资路演 (硬科技 + 潮玩)",      "追觅基金 + 招商银行 + 长江证券 + 金浦", "30+"],
            ["#4", "T+7 月",                    "设计与创意 (国潮 + 数字潮玩)",     "上海交大设计学院 + 上海市科企联",         "30+"],
            ["#5", "T+9 月",                    "内容 IP 与 Z 世代",              "中国百货协会潮玩次元专委 + 中动漫",       "30+"],
            ["#6", "T+11 月",                  "政策补贴与小镇规划",                "闵行科协 + 复旦住房政策研究中心",           "30+"],
        ])
    add_body(s, [
        "每场目标:30+ 产业客户 / 5+ 投资人 / 3+ 政府部门 / 100+ 媒体声量 → 累计触达 ≥ 200 家目标企业",
    ], top=4.6, height=0.8, font_size=14)
    add_footer(s)

    # ============== P16 Phase 3 五项挂牌 ==============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P16 · Phase 3 · 品牌与活动 (二)：5 项挂牌 (5 × 30 万 = 150 万一次性激励)")
    add_table(s, 1.2,
        ["序号", "挂牌名称", "出牌方", "对元谷的价值"],
        [
            ["①", "AI 潮玩产业基地",                "中国动漫集团",            "牌照即招商, AI+潮玩双赛道核心抓手"],
            ["②", "潮玩次元商业专委会",              "中国百货商业协会",         "聚集潮玩零售生态 + 政府对话渠道"],
            ["③", "复旦大学住房政策研究中心 · 元谷分中心", "复旦大学住房政策研究中心",  "学术背书 + 政策研究 + 高净值人脉"],
            ["④", "上海市科技企业联合会 · 元谷产业基地", "上海市科技企业联合会",     "上海科技企业生态导流 + 政府补贴申报"],
            ["⑤", "福布斯产业影响力奖 · 元谷专场",     "福布斯",                  "国际品牌势能 + 年度评选 IP"],
        ])
    add_body(s, [
        "建议森马每项挂牌支付 ¥30 万一次性激励 → 5 项合计 ¥150 万",
        "12 个月内 5 项全数挂牌 (T+3 月: ①+②; T+6 月: ④; T+9 月: ③; T+12 月: ⑤)",
    ], top=4.3, height=1.6, font_size=14)
    add_footer(s)

    # ============== P17 Phase 3 品牌 IP 化 ==============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P17 · Phase 3 · 品牌与活动 (三)：『活动 IP 化』年度日历")
    add_body(s, [
        "【一年八节奏】(以 12 个月为周期)",
        "  · 月度: 1 场产业沙龙 (年 6-12 场)",
        "  · 季度: 1 场北欧创新国际会客厅外事接待",
        "  · 半年: 1 场福布斯 元谷专场榜单发布",
        "  · 年度: 1 届全国潮玩设计大赛 (与上海交大 + 中动漫联办)",
        "  · 年度: 1 届 AI 潮玩产业大会 (合资公司主办, 邀请 1,000+ 产业人)",
        "【与 5/22 峰会的承接】每年 5 月固定承接『AI 商业化峰会』作为元谷品牌主场",
        "【可量化产出】",
        "  · 媒体声量年度 ≥ 5 亿次曝光",
        "  · 政府补贴年度 ≥ 100 万元 (申报路径)",
        "  · 赞助 / 票务年度 ≥ 200 万元",
    ], font_size=15)
    add_footer(s)

    # ============== P18 Phase 4 商业条款 (一) 月费 ==============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P18 · Phase 4 · 商业条款 (一)：基础月费 (1–3 人配置阶梯)")
    add_table(s, 1.2,
        ["配置", "人员", "月度刚性成本 (元)", "建议月费 (元/月)", "毛利率"],
        [
            ["1 人轻配 (招商单兵)",           "产业招商经理 × 1 + CSO 顾问 (折半)",                        "45,000",  "60,000",  "25%"],
            ["★ 2 人推荐配 (招商 + 活动)",   "招商经理 + 国际合作&活动策划 + CSO 顾问 + 仲量联行接口",     "95,000",  "120,000", "21%"],
            ["3 人重配 (+基金/政府)",         "招商 + 活动 + 基金投后&政府关系 + CSO 顾问 + 数据 + 行政",   "140,000", "180,000", "22%"],
        ], header_color=PRIMARY)
    add_body(s, [
        "★ 推荐方案: 2 人配置 + 12 万元/月 → 24 个月合计基础费用 ≈ 288 万元",
        "★ 1 万 ≤ 月毛利垫 ≤ 2.5 万: 用于差旅 + 接待 + 物料 + 数据运维, 留出适度安全边际",
        "成本反算依据: 详见 Excel Sheet '03 基础费用构成'",
    ], top=4.0, height=2.3, font_size=14)
    add_footer(s)

    # ============== P19 Phase 4 商业条款 (二) 佣金 ==============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P19 · Phase 4 · 商业条款 (二)：招商佣金阶梯 (1.5–2 个月年租金)")
    add_table(s, 1.2,
        ["面积档位", "日租金 (元/㎡/天)", "年租金 (元/㎡)", "佣金月数", "佣金 (元/㎡)"],
        [
            ["≤ 2,000㎡ 小型",        "2.0", "730", "1.5 个月", "91"],
            ["2,001–5,000㎡ 中型",    "2.2", "803", "1.75 个月", "117"],
            ["> 5,000㎡ 头部 / 央企", "2.5", "913", "2.0 个月",  "152"],
        ])
    add_body(s, [
        "★ 满租 2 万㎡ × 按 1.75 个月加权平均 → 佣金累计 ≈ 234 万元 (24 个月)",
        "★ 追觅基金返投落地客户额外加成 +0.25 个月 (即上限 2.25 个月)",
        "★ 客户名单提交后归属保护期 24 个月; 甲方不得绕开合资公司签约",
        "★ 与 v1.1 全盘方案的差异: 全盘方案佣金 1.0/1.5/2.5 个月偏低端; 本专项采用 1.5/1.75/2.0 中段、更激励招商团队",
    ], top=3.8, height=2.4, font_size=14)
    add_footer(s)

    # ============== P20 Phase 4 商业条款 (三) 挂牌+沙龙 ==============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P20 · Phase 4 · 商业条款 (三)：挂牌奖励 + 沙龙费 + KPI 对赌")
    add_body(s, [
        "【挂牌奖励 (一次性, 5 项)】",
        "  · 每项 ¥30 万 × 5 项 = 总额 ¥150 万",
        "  · 结算节奏: 挂牌正式公告之日起 30 日内支付",
        "【沙龙执行费 (6 场)】",
        "  · 每场 ¥5 万执行包 (场地 + 物料 + 嘉宾接待 + 传播) × 6 场 = ¥30 万",
        "  · 赞助 / 票务 / 政府补贴 → 合资公司收取, 净利 30/70 (甲方/丙方) 分润",
        "  · 单场净利预估 ¥6-8 万 → 24 个月 6 场净利 ¥36-48 万, 丙方 70% 分润 ¥25-34 万",
        "【KPI 对赌】",
        "  · T+6 月签约 ≥ 4,000㎡ / T+12 月 ≥ 10,000㎡ (4# 楼满租) / T+24 月 ≥ 18,000㎡ (90% 满租视为达标)",
        "  · 沙龙: 每场 ≥ 30 目标产业客户 (现场签到 + 闭环管理)",
        "  · 挂牌: 12 个月内 ≥ 3 项 / 24 个月内 ≥ 5 项",
        "  · 未达标的, 当期月费保留 50%, 待达成后补足",
    ], font_size=14)
    add_footer(s)

    # ============== P21 Phase 5 落地推进 ==============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P21 · Phase 5 · 落地推进 (一)：30 / 60 / 90 / 180 / 365 / 730 天行动表")
    add_table(s, 1.2,
        ["节点", "关键动作", "里程碑产出"],
        [
            ["30 天 (5/22 当周内)",  "MOU 签订 + 2-3 人战队就位 + 仲量联行数据接入 + 5/22 峰会借势完成",                   "首批 20 家意向客户清单"],
            ["60 天",                "沙龙 #1 (AI+潮玩) 完成 + AI 潮玩产业基地挂牌筹备",                                   "签约首单 + 媒体头部曝光"],
            ["90 天",                "AI 潮玩产业基地正式挂牌 + 潮玩次元专委挂牌 + 4,000㎡ 签约",                              "L1 牌照客户 3 家入驻"],
            ["180 天",               "沙龙过半 (3 场) + 上海市科企联挂牌 + 8,000㎡ 签约",                                     "L2 资本招商客户落地"],
            ["365 天",               "6 场沙龙完成 + 5 项挂牌全部落地 + 4# 楼 1 万方满租",                                   "首年报告 + 续约谈判"],
            ["730 天",               "5# 楼 1 万方满租 + 沙龙年度 IP 化 + 福布斯榜单元谷专场上线",                            "★ 2 万方目标完成"],
        ])
    add_footer(s)

    # ============== P22 Phase 5 团队与资源到位 ==============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P22 · Phase 5 · 落地推进 (二)：团队配置 + 资源接入路径")
    add_body(s, [
        "【团队配置 — 2 人核心 + CSO + 兼职专家池】",
        "  · 产业招商经理 (全职): 22K 底薪 + 10K 绩效 + 招商提成",
        "  · 国际合作 & 活动策划 (全职): 20K 底薪 + 8K 绩效 + 沙龙分润",
        "  · CSO (胡教授本人): 每周 ≥ 2 个工作日 + 重大节点全勤, 月顾问费 25K",
        "  · 兼职专家池 (项目制): 财税 / 法律 / 知识产权 / 政府补贴顾问",
        "【资源接入路径 — 90 天三轨并行】",
        "  · 数据轨: 90 日内接入仲量联行爬楼大数据 → 200 家目标客户清单",
        "  · 资本轨: 90 日内完成追觅基金 + AI 腾讯生态合作框架协议",
        "  · 牌照轨: 90 日内挂牌『AI 潮玩产业基地』+『潮玩次元商业专委会』",
        "【与森马的协同接口】",
        "  · 月度双周会 (CSO 主持) / 季度董事会汇报 / 年度战略复盘",
    ], font_size=14)
    add_footer(s)

    # ============== P23 投决建议 (一) 财务总账 ==============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P23 · 投决建议 (一)：24 个月双向账本")
    add_table(s, 1.2,
        ["收入类别 (向胡教授团队结算)", "T+1 年 (元)", "T+2 年 (元)", "24 个月合计 (元)", "说明"],
        [
            ["基础月费 (12 万 × 12)",         "1,440,000", "1,440,000", "2,880,000", "2 人配置 × 24 月"],
            ["招商佣金",                     "1,170,000", "1,170,000", "2,340,000", "按 1.75 个月平均、每年 1 万方"],
            ["挂牌奖励 (5 × 30 万)",         "1,200,000",   "300,000", "1,500,000", "4 项前置 + 1 项 T+2 年"],
            ["沙龙净分润 (70% × 30 万均值)",  "210,000",     "210,000",   "420,000", "6 场沙龙 / 年"],
            ["合计 (胡教授团队收入)",          "4,020,000", "3,120,000", "7,140,000", "≈ 714 万元 (24 个月)"],
        ])
    add_table(s, 4.5,
        ["对甲方 (森马) 的贡献", "T+1 年 (元)", "T+2 年 (元)", "24 个月合计 (元)", "说明"],
        [
            ["年化租金 (满租后)",          "8,030,000", "16,060,000", "24,090,000", "T+1 平均半租 / T+2 满租"],
            ["物业费 (10 元/㎡/月)",        "1,200,000",   "2,400,000",   "3,600,000", "随入驻渐进"],
            ["合计 (甲方直接收入)",        "9,230,000", "18,460,000", "27,690,000", "≈ 2,769 万元"],
        ], header_color=ACCENT)
    add_footer(s)

    # ============== P24 投决建议 (二) ROI 总结 ==============
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P24 · 投决建议 (二)：ROI 与决策点")
    add_kpi_card(s, 0.6, 1.3, 3.0, 1.5, "甲方付出 (24 月)", "714 万", "月费+佣金+挂牌+沙龙")
    add_kpi_card(s, 3.7, 1.3, 3.0, 1.5, "甲方收入 (24 月)", "2,769 万", "租金 + 物业")
    add_kpi_card(s, 6.8, 1.3, 3.0, 1.5, "ROI", "1 : 3.9", "投入产出比")
    add_kpi_card(s, 9.9, 1.3, 3.0, 1.5, "资产增值", "不可量化", "TOD 板块整体溢价")
    add_body(s, [
        "★ 决策三问 (建议森马本月内确定):",
        "  ① 是否接受 2 万方专项独家招商委托 (24 个月)? — 建议: 是",
        "  ② 是否接受『基础月费 12 万/月 + 1.5-2 个月佣金 + 5 项挂牌 + 6 场沙龙』条款? — 建议: 是",
        "  ③ 是否同意 30 天内启动 (5/22 峰会借势)? — 建议: 是 (错过 5/22 后该资源释放至下一年度)",
        "",
        "★ 风险预案: 若 24 个月未达满租 90% (即 ≥ 1.8 万㎡), 月费保留 50% 待达成补足 → 甲方下行可控",
        "★ 上行空间: 若 18 个月即满租, 节余成本可投入 5.2 万㎡商业部分的二期招商委托",
    ], top=3.2, height=4.0, font_size=14)
    add_footer(s)

    # ============== P25 结尾 ==============
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = PRIMARY; bg.line.fill.background()
    box = s.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(11.3), Inches(3.5)).text_frame
    box.word_wrap = True
    p1 = box.paragraphs[0]; p1.text = "AI 潮玩 · 双牌照 · 6 场沙龙 · 2 万方"
    p1.font.size = Pt(40); p1.font.bold = True; p1.font.color.rgb = WHITE
    p2 = box.add_paragraph(); p2.text = "把 5/22 峰会的高净值人群,"
    p2.font.size = Pt(28); p2.font.color.rgb = WHITE
    p3 = box.add_paragraph(); p3.text = "30 天内,变成元谷的客户。"
    p3.font.size = Pt(28); p3.font.bold = True; p3.font.color.rgb = ACCENT
    p4 = box.add_paragraph(); p4.text = ""
    p5 = box.add_paragraph(); p5.text = "胡教授团队 × 森马集团  ·  4#+5# 楼 2 万方专项招商运营方案"
    p5.font.size = Pt(16); p5.font.color.rgb = RGBColor(0xCF, 0xD8, 0xE3)

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
