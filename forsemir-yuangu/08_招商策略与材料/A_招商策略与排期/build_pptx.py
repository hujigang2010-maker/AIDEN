"""元谷 2 万方招商策略与排期主汇报 PPT.

基于 2026/6 会议纪要的战略调整版本:
  - 战略主轴: 二次元 → IP + AI 双轨
  - 硬节点 1: 2026/9/30 完成 2,000㎡ 签约 (T+100 天)
  - 硬节点 2: 2027/5/1 项目开业 (T+314 天, 需 50%+ 签约)
  - 行业转化率 110:1, 大客户优先
  - 深度绑定闵行区政府 + 腾讯算力 + CVC 基金 + 欧洲驻沪机构 + 福布斯
"""
from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

OUT = Path(__file__).with_name("元谷2万方招商策略与排期.pptx")

PRIMARY = RGBColor(0x14, 0x2C, 0x5E)
ACCENT = RGBColor(0xF2, 0x7E, 0x2D)
GOLD = RGBColor(0xC8, 0x99, 0x3D)
RED = RGBColor(0xC0, 0x39, 0x2B)
LIGHT = RGBColor(0xF5, 0xF7, 0xFA)
TEXT = RGBColor(0x1F, 0x2A, 0x44)
MUTED = RGBColor(0x55, 0x60, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.95))
    bar.fill.solid(); bar.fill.fore_color.rgb = PRIMARY; bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.4); tf.margin_top = Inches(0.15)
    tf.text = title
    p = tf.paragraphs[0]; p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = WHITE
    if subtitle:
        sub = tf.add_paragraph(); sub.text = subtitle
        sub.font.size = Pt(12); sub.font.color.rgb = RGBColor(0xCF, 0xD8, 0xE3)


def add_body(slide, bullets, left=0.6, top=1.25, width=12.1, height=5.7, font_size=14):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame; tf.word_wrap = True
    for i, line in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        depth = 0; text = line
        if text.startswith("      "):
            depth = 3; text = text.lstrip()
        elif text.startswith("    "):
            depth = 2; text = text.lstrip()
        elif text.startswith("  "):
            depth = 1; text = text.lstrip()
        para.text = text
        para.level = depth
        para.font.size = Pt(font_size if depth == 0 else font_size - 1)
        para.font.color.rgb = TEXT if depth == 0 else MUTED
        para.font.bold = depth == 0 and (text.startswith("【") or text.startswith("★") or text.startswith("【硬节点"))
        para.space_after = Pt(4)


def add_footer(slide, text="元谷 2 万方招商策略与排期 v1.0 · 胡教授团队"):
    box = slide.shapes.add_textbox(Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.4))
    tf = box.text_frame; tf.text = text
    p = tf.paragraphs[0]; p.font.size = Pt(10); p.font.color.rgb = MUTED


def add_table(slide, top_inch, headers, rows, left=0.5, width=12.3, header_color=None, row_h=0.42):
    cols = len(headers); n = len(rows) + 1
    header_color = header_color or PRIMARY
    table_shape = slide.shapes.add_table(n, cols, Inches(left), Inches(top_inch), Inches(width), Inches(0.5 + row_h * len(rows)))
    table = table_shape.table
    for j, h in enumerate(headers):
        cell = table.cell(0, j); cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = WHITE; r.font.bold = True; r.font.size = Pt(11)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j); cell.text = val
            cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT if i % 2 else WHITE
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.color.rgb = TEXT; r.font.size = Pt(10)


def add_kpi(slide, l, t, w, h, label, value, sub, color=PRIMARY):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT; box.line.color.rgb = color
    tf = box.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.1)
    tf.text = label
    p = tf.paragraphs[0]; p.font.size = Pt(10); p.font.color.rgb = MUTED
    p2 = tf.add_paragraph(); p2.text = value
    p2.font.size = Pt(18); p2.font.color.rgb = color; p2.font.bold = True
    p3 = tf.add_paragraph(); p3.text = sub
    p3.font.size = Pt(9); p3.font.color.rgb = MUTED


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ====== P1 封面 ======
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = PRIMARY; bg.line.fill.background()
    acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(3.0), Inches(0.3), Inches(2.4))
    acc.fill.solid(); acc.fill.fore_color.rgb = ACCENT; acc.line.fill.background()
    tf = s.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(11.5), Inches(3.5)).text_frame
    tf.text = "元谷项目 4#+5# 楼 2 万方"
    tf.paragraphs[0].font.size = Pt(48); tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = WHITE
    p = tf.add_paragraph(); p.text = "IP + AI 双轨招商策略与排期"
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = ACCENT
    p2 = tf.add_paragraph(); p2.text = ""
    p3 = tf.add_paragraph(); p3.text = "硬节点 ① 2026/9/30 — 2,000㎡ 签约"
    p3.font.size = Pt(18); p3.font.color.rgb = WHITE
    p4 = tf.add_paragraph(); p4.text = "硬节点 ② 2027/5/1 — 项目开业 (50%+ 签约)"
    p4.font.size = Pt(18); p4.font.color.rgb = WHITE
    p5 = tf.add_paragraph(); p5.text = ""
    p6 = tf.add_paragraph(); p6.text = "胡教授团队  ×  森马集团  ·  2026/6 提案 v1.0"
    p6.font.size = Pt(14); p6.font.color.rgb = RGBColor(0xCF, 0xD8, 0xE3)

    # ====== P2 执行摘要 ======
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P2 · 执行摘要", "30 秒读懂本方案 — 战略转向 IP+AI 双轨")
    add_body(s, [
        "★ 战略调整: 主轴由'纯二次元'升级为 'IP + AI 双轨' (AI 拿政策/基金/税收, IP 兼顾威总文化导向)",
        "★ 两个硬节点 (按 2026/6/22 起算):",
        "  · 2026/9/30 → 2,000㎡ 签约 (T+100 天, 可含直播基地/共享设计中心擦边球)",
        "  · 2027/5/1 → 项目开业 (T+314 天, 2 万方需 50%+ 签约 + 消费氛围成型)",
        "★ 客户策略: 优先招 2,000-20,000㎡ 大客户; 行业转化率 110:1 (招 1 单要见 110 家)",
        "★ 政府绑定: 闵行区科委 + 商务委 + 街道 三轨并行, 加 IP+AI 主线",
        "★ 资源拉通: 腾讯算力补贴 + CVC 创投基金 + 欧洲驻沪机构(德 / 以 / 北欧) + 福布斯",
        "★ 团队建议: 2 人驻场 + CSO 顾问 (轻骑兵, 月度成本可控)",
        "★ 商业条款: 基础月费 12 万/月 + 招商佣金 1.5-2 个月年租金 + 5 项挂牌 30 万/项 + 6 场活动 5 万/场",
        "★ 24 月 ROI: 甲方付出 ≈ 714 万元 / 甲方收入 ≈ 2,769 万元 / 投入产出 1:3.9",
    ], font_size=14)
    add_footer(s)

    # ====== P3 战略转向(核心) ======
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P3 · 战略转向 (核心) — 从'纯二次元' → 'IP + AI 双轨'")
    add_table(s, 1.2,
        ["维度", "原方向 (纯二次元)", "新方向 (IP + AI 双轨)"],
        [
            ["招商面",         "二次元企业 200-500㎡ 小户, 招满需 400-800 家",                 "AI/IP 中大型 1,000-5,000㎡, 招满仅需 30-60 家"],
            ["支付能力",       "二次元企业规模小, 单家年租金 < 30 万",                          "AI/IP 企业真金白银, 单家年租金 100-500 万"],
            ["政府政策",       "二次元拿政策困难, 创新基金不投",                                "AI 享高新认定 + 专精特新 + 创新券 + 闵行专项"],
            ["与森马协同",     "二次元与森马商业难协同",                                       "IP+AI 与森马时尚科技定位高度契合"],
            ["先例",          "上海无 2 万方纯二次元成功项目",                                "杨浦区 1 万方 'AI+IP' 项目已通过先例"],
            ["招商速度",       "脱离政府主线, 速度慢 50%+",                                    "深度绑定闵行区科委/商务委, 速度翻倍"],
        ], row_h=0.45)
    add_body(s, [
        "★ 结论: 主轴改 AI 同时保留 IP 概念 (涵盖视频/漫画/直播/动漫全链条), 既能拿政府资源, 又兼顾威总战略文化导向",
    ], top=5.2, height=0.7, font_size=13)
    add_footer(s)

    # ====== P4 IP+AI 双轨产业地图 ======
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P4 · IP + AI 双轨产业地图 — 4# 楼 + 5# 楼 5F+ 各自定位")
    add_table(s, 1.2,
        ["楼栋 (5F+)", "面积", "主轴", "目标客户类型", "代表企业画像"],
        [
            ["4# 楼", "≈ 1 万㎡", "★ AI 主轴 (国际化层)", "AI 设计 / AIGC / AI 内容 / 国际 IP / AI Agent / 大模型应用", "字节 / MiniMax / 月之暗面 / 360AI / 智谱 / 海外 AI 公司"],
            ["5# 楼", "≈ 1 万㎡", "★ IP 主轴 (产业总部层)", "潮玩 / 动漫 / 数字 IP / 内容 / 直播 / 衍生品", "中动漫旗下 / 上海美影厂周边 / 头部 IP 运营商 / B 站生态"],
        ], row_h=0.6)
    add_body(s, [
        "★ 4# 楼 同栋协同: 4F 直播中心 + 5F AI 共享设计中心 + 5F AI 共享打样&DIY 中心",
        "★ 5# 楼 同栋协同: 5F 潮玩产业展厅 + 1-4F 动漫书店&休闲娱乐",
        "★ 4# 楼 价高: 建议 2.2-2.5 元/㎡/天 (主轴 AI 享区域品牌溢价)",
        "★ 5# 楼 价稳: 建议 2.0-2.2 元/㎡/天 (主轴 IP 产业, 与同栋娱乐配套形成生态)",
        "★ 双轨融合: '杨浦区 AI+IP 1 万方先例' 在元谷 4#/5# 楼形成 2 万方升级版",
    ], top=3.8, height=3.0, font_size=13)
    add_footer(s)

    # ====== P5 两个硬节点 ======
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P5 · 两个硬节点 — 倒推全部排期", "倒计时, 不可妥协")
    add_kpi(s, 0.6, 1.3, 6.0, 1.6, "节点一", "2026/9/30 · T+100 天", "完成 2,000㎡ 签约 (硬指标)", color=RED)
    add_kpi(s, 6.7, 1.3, 6.0, 1.6, "节点二", "2027/5/1 · T+314 天", "项目开业, 2 万方 50%+ 签约 + 消费氛围", color=ACCENT)
    add_table(s, 3.2,
        ["阶段", "起止", "招商目标", "累计 / 阶段贡献"],
        [
            ["阶段 0 准备",    "2026/6/22 - 2026/7/15",   "0",                                  "团队就位 + 资源接入 + 物料 (T+23 天)"],
            ["阶段 1 抢节点", "2026/7/15 - 2026/9/30",   "★ 2,000㎡ (含擦边球)",               "硬节点一, 100% 必达"],
            ["阶段 2 加速",    "2026/10/1 - 2026/12/31",  "+ 3,000㎡ → 累计 5,000㎡",           "Q4 攻坚"],
            ["阶段 3 开业",   "2027/1/1 - 2027/5/1",     "+ 5,000㎡ → 累计 10,000㎡ (50%)",   "硬节点二, 项目开业"],
            ["阶段 4 满租",   "2027/5/2 - 2028/5/1",     "+ 8,000-10,000㎡ → 累计 90%+",      "全周期满租"],
        ], row_h=0.4)
    add_footer(s)

    # ====== P6 110:1 与大客户优先 ======
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P6 · 行业转化率 110:1 — 反推工作量与策略")
    add_table(s, 1.2,
        ["转化率历史", "时期", "原因"],
        [
            ["45 : 1",    "2015-2018", "经济上行, 办公需求旺盛"],
            ["68 : 1",    "2019-2022", "首轮经济调整, 需求回落"],
            ["★ 110 : 1", "2023-2026 (当前)", "AI 替代办公需求 + 经济下行 + 企业压成本"],
        ], row_h=0.45)
    add_body(s, [
        "【两种打法的工作量对比】",
        "  · 打小户 (200㎡ 二次元): 招满 2 万方 = 100 单 × 110 = 触达 11,000 家潜客 (3 人团队需 24+ 月)",
        "  · 打大户 (1,000㎡ AI/IP): 招满 2 万方 = 20 单 × 110 = 触达 2,200 家潜客 (3 人团队需 8-10 月)",
        "  → 效率差 5 倍, 节点风险差 10 倍",
        "",
        "【对应的运营策略】",
        "  · 重点行业聚焦: AI / IP 头部 + 上下游, 不广撒网",
        "  · 大客户'特批通道': 入驻 ≥2,000㎡ 客户享受 装修补贴 + 政府对接 + 服务中心免费 1 年",
        "  · 招商节奏: 月度 KPI 触达 500 家 + 深谈 30 家 + 成交 4-5 家",
        "  · 单笔现金流参考: 一家 265㎡ 客户 = 押 2 月 + 预付 3 月 = ≈ 60-70 万元入账",
    ], top=3.4, height=3.5, font_size=13)
    add_footer(s)

    # ====== P7 总体排期总览(甘特表) ======
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P7 · 总体排期甘特概览 (2026/6 - 2027/5 关键 11 个月)")
    add_table(s, 1.2,
        ["关键工作流", "6月", "7月", "8月", "9月", "10月", "11月", "12月", "1月", "2月", "3月", "4月", "5月"],
        [
            ["合同&战队就位",  "█ ", "█ ", "",  "",  "",  "",  "",  "",  "",  "",  "",  ""],
            ["资源&牌照接入", "█ ", "█ ", "█ ", "",  "",  "",  "",  "",  "",  "",  "",  ""],
            ["招商物料制作",  "█ ", "█ ", "█ ", "█ ", "",  "",  "",  "",  "",  "",  "",  ""],
            ["招商执行",      "",  "█ ", "█ ", "█★", "█ ", "█ ", "█ ", "█ ", "█ ", "█ ", "█ ", "█★"],
            ["产业沙龙 6 场", "",  "█ ", "",  "█ ", "",  "█ ", "",  "█ ", "",  "█ ", "",  "█ "],
            ["挂牌 5 项",     "",  "",  "█ ", "█ ", "",  "█ ", "",  "",  "█ ", "",  "",  "█ "],
            ["开业筹备",      "",  "",  "",  "",  "",  "",  "█ ", "█ ", "█ ", "█ ", "█ ", "█★"],
        ], row_h=0.4)
    add_body(s, [
        "★ 9/30 (硬节点一) 与 5/1 (硬节点二) 标红;  '█' = 工作进行中",
    ], top=5.2, height=0.7, font_size=12)
    add_footer(s)

    # ====== P8 阶段 0: 准备期 ======
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P8 · 阶段 0 (准备期) · 6/22 - 7/15 · T+0 ~ T+23 天")
    add_table(s, 1.2,
        ["工作项", "森马侧 (责任人 + 工作)", "胡教授团队 (责任人 + 工作)", "预计费用 (元)"],
        [
            ["合同与战队",      "★ 发哥 / 周志超 与威总对齐合同流程 + 资金比例", "★ 胡教授 + 招商经理候选人面试 + 法务审稿", "—"],
            ["资源接入",        "森马提供品宣物料 + 楼盘 CAD/效果图",            "★ 仲量联行爬楼数据接口接入 + 腾讯 + CVC 联动",  "数据接口 ¥2.6 万 (已购)"],
            ["政府对接",        "★ 危总安排闵行区科委/商务委首次拜访",           "胡教授对接复旦住房政策中心 + 上海市科企联",   "拜访礼物物料 ¥1 万"],
            ["招商物料 1.0",   "提供产品技术参数 + 楼层平面",                  "★ 招商手册 + 装修手册 + 百问话术 (本方案已交)", "设计公司外包 ¥3-5 万"],
            ["新闻发布",        "森马总部 PR 对接 + 主流媒体邀请",               "★ 胡教授对接福布斯采访 + 闵行区官网/新华网",    "PR 发布 ¥0 (媒体资源)"],
        ], row_h=0.55)
    add_footer(s)

    # ====== P9 阶段 1: 抢硬节点一 ======
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P9 · 阶段 1 (抢节点) · 7/15 - 9/30 · T+23 ~ T+100 天", "★ 硬节点一: 2,000㎡ 签约 (含擦边球)")
    add_table(s, 1.2,
        ["工作项", "森马侧", "胡教授团队", "预计费用 (元)", "对节点贡献"],
        [
            ["大客户突击 (≥ 1,000㎡)",     "★ 发哥协调签约价格 + 商务条款",     "★ L1+L2 漏斗执行: 央企 + 中型 AI/IP",          "招商交通 ¥3 万",  "签约 1-2 家 ≈ 1,000-2,000㎡"],
            ["直播基地落地 (擦边球)",       "提供 4# 楼 4F 场地 + 装修",         "★ 联动绮丽少女 + 元谷直播 IP 招标",            "硬件设备 ¥30-50 万",  "计入 ≈ 800-1,000㎡"],
            ["AI 共享设计中心 (擦边球)",    "提供 4# 楼 5F 场地",               "★ 与腾讯/AI 腾讯生态联合挂牌",                "工位 + AI 设备 ¥40-60 万",  "计入 ≈ 500-800㎡"],
            ["首批沙龙 (#1 AI+潮玩)",      "森马品牌联合宣传",                  "★ 联动腾讯 + 中动漫, 借势 5/22 峰会余热",     "执行包 ¥5 万",     "曝光 + 锁定 5-10 客户"],
            ["首批挂牌 (① + ②)",           "森马挂牌仪式 + 楼面物料",            "★ AI 潮玩产业基地 + 潮玩次元商业专委会",       "挂牌仪式 ¥10 万",    "牌照即招商"],
            ["签约率追踪",                  "森马商务部 + 财务对账",              "★ 周报机制, 数据透明",                       "—",              "—"],
        ], row_h=0.55)
    add_footer(s)

    # ====== P10 阶段 2 加速 ======
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P10 · 阶段 2 (加速期) · 10/1 - 12/31 · 累计 5,000㎡ 目标")
    add_table(s, 1.2,
        ["工作项", "森马侧", "胡教授团队", "预计费用 (元)", "对累计贡献"],
        [
            ["L3 爬楼大数据全开",          "森马 IT 对接数据接口",                "★ 200 家精准客户清单 + 上门拜访",              "差旅 ¥5 万/月",   "签约 ≈ 1,500㎡"],
            ["L2 资本招商 (CVC 落地)",     "投资协议法务支持",                    "★ 追觅基金 + 腾讯 + CVC 落地 ≥ 3 家 AI 公司",  "投资基金推介 ¥3 万",  "签约 ≈ 1,500㎡"],
            ["沙龙 #2 出海 + #3 投融资",    "森马联合宣传",                       "★ 联动北欧/福布斯 + 招行/长江/金浦/铂帝",     "执行包 ¥10 万 (2 场)", "锁定 ≈ 600㎡"],
            ["挂牌 ③ + ④",                "森马挂牌仪式",                       "★ 上海市科企联 + 复旦住房政策中心",            "挂牌仪式 ¥5 万",   "学术与政府背书"],
            ["开业筹备启动",               "★ 森马商业部主导开业策划",            "胡教授团队配合产业入驻动线设计",                "—",              "—"],
        ], row_h=0.55)
    add_footer(s)

    # ====== P11 阶段 3 开业 ======
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P11 · 阶段 3 (开业期) · 2027/1/1 - 5/1 · 累计 10,000㎡ + 开业")
    add_table(s, 1.2,
        ["工作项", "森马侧", "胡教授团队", "预计费用 (元)", "对累计贡献"],
        [
            ["大客户决战 (≥ 2,000㎡)",     "★ 威总亲自背书 + 大客户接待",         "★ 锁定头部央企 / 行业旗舰 1-2 家",            "高端接待 ¥10 万",    "签约 ≈ 2,000-4,000㎡"],
            ["沙龙 #4 + #5 + #6",          "森马联合宣传",                       "★ 上海交大 + 中百协 + 闵行科协联办",          "执行包 ¥15 万 (3 场)", "锁定 ≈ 800㎡"],
            ["挂牌 ⑤ 福布斯专场",          "森马高端晚宴",                       "★ 福布斯榜单元谷专场发布",                    "挂牌+晚宴 ¥30 万",   "国际品牌势能"],
            ["开业典礼 5/1",               "★ 森马集团董事长出席 + 政府高层邀请",  "★ 胡教授+复旦+科企联+欧洲机构出席",            "开业典礼 ¥80-100 万",  "项目品牌定型"],
            ["二期 5.2 万方招商启动 (预热)", "★ 森马商业部承接",                   "胡教授团队提供产业地图 + 客户管道",             "—",              "为下个阶段铺垫"],
        ], row_h=0.55)
    add_footer(s)

    # ====== P12 阶段 4 满租 ======
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P12 · 阶段 4 (满租期) · 2027/5/2 - 2028/5/1 · 累计 90%+")
    add_body(s, [
        "★ 阶段目标: 4#/5# 楼共 2 万方达到 90%+ 满租 (= ≥ 18,000㎡)",
        "★ 核心动作:",
        "  · 沙龙 IP 化年度化 (年度 12-15 场, 单场 ≥ 30 客户)",
        "  · 5 项挂牌续约 (每项 50% 续挂费 = 15 万/项)",
        "  · 福布斯榜单年度评选发布 (元谷主场)",
        "  · 整合 5/22 'AI 商业化峰会' 作为元谷品牌固定 IP",
        "  · 招商对象延伸至 5.2 万方商业部分 (二期承接)",
        "★ 财务模型 (第二年):",
        "  · 月费 13 万/月 × 12 = 156 万元 (可适度涨价)",
        "  · 招商佣金 (+ 8,000㎡ 增量): 1.75 个月 × 803 元/㎡ × 8,000 ÷ 12 ≈ 94 万元",
        "  · 沙龙 #7-#12 (6 场): 净分润 70% × 31 万 = 21.7 万元",
        "  · 挂牌续挂奖励: 5 × 15 万 = 75 万元",
        "  · 超额奖励 (满租率 ≥ 95% 一次性): 100 万元",
        "  · 第二年合计: ≈ 446 万元 (满租场景)",
        "★ 24 月双向账本: 甲方付出 ≈ 714 万 / 甲方收入 ≈ 2,769 万 / ROI 1:3.88",
    ], font_size=13)
    add_footer(s)

    # ====== P13 资源拉通(乙方提供) ======
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P13 · 资源拉通 (胡教授团队提供, 直接转化为招商工具)")
    add_table(s, 1.2,
        ["资源", "类型", "落地路径", "对节点贡献"],
        [
            ["腾讯算力补贴", "AI 公司专享", "新注册 AI 公司 3 个月合同免费 / 算力 85 折", "★ 锁定 ≥ 5 家 AI 公司"],
            ["CVC 创业投资基金", "资本招商", "天使-A 轮投资, 入驻元谷 AI/IP 企业优先", "★ 锁定 ≥ 3 家"],
            ["闵行区科委 / 商务委 / 街道", "政策资源", "高新认定 / 专精特新 / 创新券 / 闵行专项补贴", "提升招商转化率 30%"],
            ["欧洲驻沪机构 (德/以/北欧)", "国际化资源", "'科技 + 出海' 展示场景", "国际品牌势能 + 主流媒体声量"],
            ["复旦大学住房政策研究中心", "学术背书", "元谷分中心挂牌, 政策研究 + 高净值人脉", "高端客户对话渠道"],
            ["上海市科技企业联合会", "上海科企生态", "元谷产业基地挂牌, 上海科企导流", "上海生态对接"],
            ["福布斯产业影响力榜", "国际品牌势能", "元谷专场榜单 + 国际媒体曝光", "国际品牌势能"],
            ["仲量联行爬楼大数据", "招商弹药", "200+ 家目标客户清单 + 转化率 +30%", "★ 核心销售工具"],
        ], row_h=0.36)
    add_footer(s)

    # ====== P14 团队配置与商业条款 ======
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P14 · 团队配置与商业条款 (向森马提议)")
    add_table(s, 1.2,
        ["项", "推荐方案", "区间 / 备注"],
        [
            ["团队配置",         "★ 2 人驻场 + CSO 顾问 (胡教授每周 ≥ 2 天)",  "1-3 人阶梯可选"],
            ["基础月费",         "★ 12 万元/月",                                "6-18 万 (1/2/3 人), 推荐 2 人"],
            ["招商佣金阶梯",     "1.5 / 1.75 / 2.0 个月年租金 (按面积)",         "≤2000㎡/2-5000㎡/>5000㎡; 返投 +0.25 月"],
            ["挂牌奖励",         "5 项 × 30 万 = 150 万 (一次性)",             "续挂 50% (15 万/项)"],
            ["沙龙执行费",       "5 万/场 × 6 场 = 30 万",                    "净利 30/70 分润 (甲/乙)"],
            ["服务期",           "24 个月 (首期)",                            "续期由乙方优先选择权"],
            ["超额奖励",         "24 月满租率 ≥ 95% → 100 万一次性",             "对赌激励"],
            ["客户归属保护期",   "24 个月 (名单提交后)",                       "甲方不得绕开签约"],
        ], row_h=0.42)
    add_footer(s)

    # ====== P15 9/30 节点深拆 ======
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P15 · 9/30 节点深拆 — 2,000㎡ 怎么来", "T+100 天必达军令状")
    add_table(s, 1.2,
        ["来源", "面积估算", "落地概率", "兜底方案"],
        [
            ["1-2 家中大型 AI/IP 客户",     "1,000-2,000㎡",   "60%",   "若未签, 直播基地 + 共享设计中心兜底"],
            ["共享直播基地落地",            "800-1,000㎡",     "★ 90%", "森马已具备意愿, 7 月启动硬件采购"],
            ["AI 共享设计中心落地",          "500-800㎡",      "★ 80%", "联动腾讯/上海交大, 8 月挂牌"],
            ["产业服务中心入驻",            "300-500㎡",      "100%",  "乙方自营业态, 100% 计入"],
            ["首批小型 AI/IP 客户 3-5 家",   "600-1,500㎡",    "50%",   "L4 沙龙带流的快速转化"],
            ["合计 (中性估算)",             "★ 3,200-5,800㎡", "—",     "实际计入 ≥ 2,000㎡ 概率 > 95%"],
        ], row_h=0.42)
    add_body(s, [
        "★ 关键: 共享直播基地 + AI 共享设计中心 是 '擦边球' 主力, 9/30 之前必须落地",
    ], top=5.3, height=0.7, font_size=13)
    add_footer(s)

    # ====== P16 双方分工 ======
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P16 · 双方分工 (森马侧 / 胡教授团队侧 / 联合)")
    add_table(s, 1.2,
        ["分工类别", "森马侧主责", "胡教授团队主责", "联合"],
        [
            ["战略与合同",     "威总 + 发哥 + 周志超",       "胡教授 (产业战略)",         "签约 + 商务条款"],
            ["产业招商",       "—",                         "★ 全权负责 (2 万方)",        "—"],
            ["商业招商",       "★ 发哥 + 周志超 (5.2 万方)", "—",                          "—"],
            ["政府对接",       "森马总部 PR",                "★ 胡教授 (闵行区主导)",      "联合拜访"],
            ["资源接入",       "森马品牌资源",                "★ 腾讯/CVC/欧洲/福布斯",     "联合背书"],
            ["挂牌活动",        "森马高层出席",                "★ 牵线 + 仪式执行",          "联合发布"],
            ["开业典礼",        "★ 森马集团董事长 + 政府高层", "胡教授 + 复旦 + 科企联",      "联合主办"],
            ["传播 / PR",       "森马商业部 + 主流媒体",        "福布斯 + 政府官网 + 新华网", "联合发声"],
        ], row_h=0.4)
    add_footer(s)

    # ====== P17 月度签约率推进表 ======
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P17 · 月度签约率推进表 (按年租金 803 元/㎡ 估算)")
    add_table(s, 1.2,
        ["月份", "本月新增 (㎡)", "累计签约 (㎡)", "累计签约率", "累计年租金 (万元)", "里程碑"],
        [
            ["2026/7",  "300",   "300",     "1.5%",  "24",   "战队就位 + 物料"],
            ["2026/8",  "800",   "1,100",   "5.5%",  "88",   "首批 L1 牌照 + 大客户接触"],
            ["2026/9",  "1,200", "2,300",   "11.5%", "185",  "★ 9/30 硬节点达成"],
            ["2026/10", "1,000", "3,300",   "16.5%", "265",  "L3 爬楼启动"],
            ["2026/11", "1,000", "4,300",   "21.5%", "346",  "沙龙 #3 投融资"],
            ["2026/12", "1,200", "5,500",   "27.5%", "442",  "Q4 攻坚收官"],
            ["2027/1",  "1,200", "6,700",   "33.5%", "538",  "新年首单"],
            ["2027/2",  "1,000", "7,700",   "38.5%", "619",  "—"],
            ["2027/3",  "1,200", "8,900",   "44.5%", "715",  "—"],
            ["2027/4",  "1,500", "10,400",  "★ 52%", "836",  "★ 5/1 节点提前一周达成"],
            ["2027/5",  "1,500", "11,900",  "59.5%", "956",  "开业典礼"],
        ], row_h=0.32)
    add_footer(s)

    # ====== P18 费用与现金流 ======
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P18 · 12 个月费用与现金流 (向森马提案口径)")
    add_table(s, 1.2,
        ["费用类别", "金额 (元)", "支付节奏", "甲方现金压力 (元/月)"],
        [
            ["团队月费 12 万 × 12",          "1,440,000", "按月预付",          "120,000"],
            ["招商佣金 (累计 11,900㎡ 估)",   "1,400,000", "起租后 30 日内",     "渐进, 高峰约 200,000"],
            ["5 项挂牌 (一次性首挂)",         "1,500,000", "挂牌 30 日内",      "分 5 次, 平均 25,000"],
            ["6 场沙龙执行 (5 万/场)",        "300,000",   "按场结算",          "约 25,000"],
            ["首场开业典礼 (5/1)",            "800,000",   "开业前 30 日内",     "—"],
            ["资源对接 / 政府关系 / 接待",     "200,000",   "实报实销",          "—"],
            ["12 月合计",                    "★ 5,640,000", "—",                "—"],
            ["月均现金压力",                 "—",          "—",                "★ ≈ 47 万/月 (含佣金高峰)"],
        ], row_h=0.4)
    add_body(s, [
        "★ 同期甲方租金累计收入 ≈ 956 万元, 物业 ≈ 90 万元 → 单 12 个月已现金回正",
    ], top=5.5, height=0.6, font_size=13)
    add_footer(s)

    # ====== P19 风险与对冲 ======
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P19 · 风险与对冲 (写进协议)")
    add_table(s, 1.2,
        ["风险", "影响", "对冲机制"],
        [
            ["9/30 节点 2,000㎡ 未达",         "森马失信 + 项目延期",         "★ 擦边球 (直播 + 设计中心) 兜底 + KPI 对赌月费保留 50%"],
            ["5/1 开业 50% 未达",              "开业氛围不足",                 "★ 阶段 4 满租期延长 + 沙龙 IP 化加速"],
            ["转化率 110:1 → 130:1 (恶化)",     "招商节奏放缓 30%",            "★ 团队扩 1 人 (改 3 人重配) + 月费上调至 18 万"],
            ["闵行区政策变化",                  "补贴/牌照不落地",              "★ 备份: 杨浦区 / 漕河泾 / 张江政策可援引"],
            ["AI 行业资本退潮",                 "AI 客户付款能力下降",          "★ 切换 IP 主轴比例 + 加快出海企业引入"],
            ["森马商业部协同延误",              "1#/6# 商业氛围不足",          "★ 5.2 万方与 2 万方独立排期 + 不相互拖累"],
        ], row_h=0.42)
    add_footer(s)

    # ====== P20 投决建议 ======
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P20 · 投决建议", "建议本周内启动")
    add_body(s, [
        "★ 三问决策 (建议森马本周内确定):",
        "  ① 是否接受 IP+AI 双轨战略调整? — 建议: 是 (杨浦区已有先例)",
        "  ② 是否接受 2 人驻场 + 12 万元月费 + 1.5-2 月佣金 + 5 项挂牌 + 6 场沙龙条款? — 建议: 是",
        "  ③ 是否同意 7 月初启动 (赶 9/30 硬节点)? — 建议: 是, 错过 7 月即风险陡增",
        "",
        "★ 一句话总结:",
        "  以 IP+AI 双轨战略, 借势 5/22 峰会余热与闵行区科委政策红利,",
        "  100 天达成 9/30 节点 2,000㎡, 314 天达成 5/1 开业 50% 签约,",
        "  730 天达成 2 万方满租, 为森马贡献年化租金 1,606 万元 + 资产估值 2.47 亿元",
        "",
        "★ 决策成本:",
        "  ─ 投入: 12 个月 ≈ 564 万元 (月费 + 佣金 + 挂牌 + 沙龙 + 开业)",
        "  ─ 收入: 同期租金 ≈ 956 万 + 物业 ≈ 90 万 = ≈ 1,046 万元",
        "  ─ 净收益: ≈ +482 万元 (12 个月单算已正), 24 个月 ROI 1:3.88",
    ], font_size=14)
    add_footer(s)

    # ====== P21 收尾页 ======
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = PRIMARY; bg.line.fill.background()
    box = s.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(11.3), Inches(3.5)).text_frame
    box.word_wrap = True
    p1 = box.paragraphs[0]; p1.text = "100 天 · 2,000 方"
    p1.font.size = Pt(56); p1.font.bold = True; p1.font.color.rgb = WHITE
    p2 = box.add_paragraph(); p2.text = "IP + AI 双轨, 借势政府, 闭环招商。"
    p2.font.size = Pt(28); p2.font.color.rgb = WHITE
    p3 = box.add_paragraph(); p3.text = ""
    p4 = box.add_paragraph(); p4.text = "胡教授团队 × 森马集团  ·  元谷 2 万方专项招商策略与排期"
    p4.font.size = Pt(16); p4.font.color.rgb = RGBColor(0xCF, 0xD8, 0xE3)

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
