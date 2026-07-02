"""Generate the Yuangu招商运营战略 PPT deck (.pptx).

The deck mirrors `PPT汇报大纲.md` so the markdown outline and the
binary deck stay in sync. Run with `python build_pptx.py`; the output
file `元谷项目招商运营战略汇报.pptx` is written next to this script.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

OUT = Path(__file__).with_name("元谷项目招商运营战略汇报.pptx")

PRIMARY = RGBColor(0x14, 0x2C, 0x5E)
ACCENT = RGBColor(0xF2, 0x7E, 0x2D)
LIGHT = RGBColor(0xF5, 0xF7, 0xFA)
TEXT = RGBColor(0x1F, 0x2A, 0x44)
MUTED = RGBColor(0x55, 0x60, 0x7A)


def add_title_bar(slide, title: str, subtitle: str | None = None) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.15)
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    if subtitle:
        sub = tf.add_paragraph()
        sub.text = subtitle
        sub.font.size = Pt(12)
        sub.font.color.rgb = RGBColor(0xCF, 0xD8, 0xE3)


def add_body(slide, bullets: list[str], left=0.6, top=1.2, width=12.1, height=5.7, font_size=18) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        depth = 0
        text = line
        if text.startswith("    "):
            depth = 2
            text = text.lstrip()
        elif text.startswith("  "):
            depth = 1
            text = text.lstrip()
        para.text = text
        para.level = depth
        para.font.size = Pt(font_size if depth == 0 else font_size - 2)
        para.font.color.rgb = TEXT if depth == 0 else MUTED
        para.font.bold = depth == 0 and text.startswith("【")
        para.space_after = Pt(6)


def add_footer(slide, text: str = "胡教授团队 × 森马集团 · 元谷项目联合运营方案") -> None:
    box = slide.shapes.add_textbox(Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.4))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED


def add_table(slide, top_inch: float, headers: list[str], rows: list[list[str]], left=0.6, width=12.1) -> None:
    cols = len(headers)
    n = len(rows) + 1
    table_shape = slide.shapes.add_table(n, cols, Inches(left), Inches(top_inch), Inches(width), Inches(0.5 + 0.45 * len(rows)))
    table = table_shape.table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                r.font.bold = True
                r.font.size = Pt(13)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if i % 2 else RGBColor(0xFF, 0xFF, 0xFF)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.color.rgb = TEXT
                    r.font.size = Pt(11)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # P1 Cover
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = PRIMARY; bg.line.fill.background()
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(3.2), Inches(0.3), Inches(1.6))
    accent.fill.solid(); accent.fill.fore_color.rgb = ACCENT; accent.line.fill.background()
    title = s.shapes.add_textbox(Inches(1.1), Inches(2.8), Inches(11.5), Inches(2.4)).text_frame
    title.text = "元谷"
    title.paragraphs[0].font.size = Pt(72); title.paragraphs[0].font.bold = True
    title.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p2 = title.add_paragraph(); p2.text = "大零号湾科技潮玩产业策源高地 · 招商运营战略"
    p2.font.size = Pt(28); p2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p3 = title.add_paragraph(); p3.text = "胡教授团队 × 森马集团 联合运营方案 v1.0"
    p3.font.size = Pt(18); p3.font.color.rgb = RGBColor(0xCF, 0xD8, 0xE3)

    # P2 双重身份
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P2 元谷的双重战略身份", "大零号湾文创融合核心区 × 上海唯一科技时尚特色小镇")
    add_body(s, [
        "【两大稀缺身份】",
        "  大零号湾文创融合核心区：闵行五大中心之一，能级比肩漕河泾、张江",
        "  上海市唯一科技时尚特色小镇：政策与品牌双重独占",
        "【建筑底盘】",
        "  总建面 22 万㎡ / 商业 5.2 万㎡ / 1500+ 商业车位",
        "  元江路单日客流 5–7 万人次 · 15 号线元江路站 TOD 项目",
        "【15 分钟覆盖】",
        "  居住人口 24 万 + 产业办公人口 12 万",
        "【一句话定位】",
        "  森马给了元谷骨架，我们要给它产业灵魂。",
    ])
    add_footer(s)

    # P3 产业判断
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P3 我们对元谷的产业判断")
    add_body(s, [
        "1#–6# 既定功能（总部 / Livehouse / 酒店 / 潮玩艺术中心 / 动漫书店 / 餐饮商务）已是优质零售底座。",
        "纯零售逻辑无法支撑“产业策源”，需补三条产业能级链：",
        "  国际链：海外 IP / 技术 → 中国首站",
        "  资本链：返投基金 → 锁定潮玩中下游产能",
        "  数据链：精准画像 → 把客流转化为产业线索",
    ])
    add_footer(s)

    # P4 三件套
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P4 三件套核心打法：基金 + 基地 + 活动")
    add_body(s, [
        "【基金】追觅科技基金 + AI 腾讯导流 → “返投落地”招商",
        "【基地】北欧创新国际会客厅平移 4# 楼 + 两大产业牌照（AI 潮玩产业基地 / 潮玩次元商业专委会）",
        "【活动】科技开放麦 + 全国潮玩设计大赛 + 福布斯榜单 → 活动即招商即媒介",
        "三件套共同构成元谷飞轮：流量 → 品牌势能 → 招商转化 → 资本助推 → 再放大流量",
    ])
    add_footer(s)

    # P5 北欧会客厅
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P5 资源盘点 (1) 北欧创新国际会客厅")
    add_body(s, [
        "已成熟运营的国际化 IP 平台，覆盖北欧设计 / ESG / 潮玩 / AI / 机器人创新生态。",
        "拟整体平移至元谷 4# 楼：",
        "  4F 直播中心 → 北欧 IP 中国首播首发地",
        "  1–3F 潮玩艺术中心 → 国际潮玩艺术家入驻 / 限定快闪",
        "价值：把 4# 楼从“网红打卡点”升级为“国际潮玩外事窗口”。",
    ])
    add_footer(s)

    # P6 福布斯 + 牌照
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P6 资源盘点 (2) 福布斯奖项 + 行业牌照")
    add_body(s, [
        "福布斯系列奖项 → 为森马总部 1# 楼提供年度品牌势能爆发节点",
        "中国百货商业协会“潮玩次元商业专委会”牌照（已确认）",
        "中国动漫集团“AI 潮玩产业基地”牌照（已确认）",
        "配合“科技时尚特色小镇”申报、闵行五大中心政策包",
        "形成「奖项 + 牌照 + 小镇」三章合一",
    ])
    add_footer(s)

    # P7 科技开放麦
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P7 资源盘点 (3) 科技开放麦 × 潮玩大赛")
    add_body(s, [
        "“科技开放麦”已运营多季，主题：AI、机器人、设计创新、出海等",
        "单场触达 300–500 高质量产业人，已具备招商转化能力",
        "元谷主场化：每月 1 场科技开放麦 + 每年 1 届全国潮玩设计大赛",
        "落地 2# 楼二次元 Livehouse / 秀场，活动现场即招商现场",
    ])
    add_footer(s)

    # P8 AI腾讯 + 仲量联行 + 追觅
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P8 资源盘点 (4) AI 腾讯 + 仲量联行 + 追觅基金")
    add_table(s, 1.3,
        ["资源", "角色", "在元谷的应用"],
        [
            ["AI 腾讯", "技术 + 内容引擎", "赋能 4# 楼 AI 共享设计中心、5#/6# 楼数字化运营"],
            ["仲量联行爬楼大数据 (¥2.6 万已购入)", "招商弹药库", "闵行 / 临港 / 漕河泾 / 张江存量企业精准画像 + 爬楼行动"],
            ["追觅科技基金", "返投资本", "1:1.5 返投比例，锁定潮玩中下游 + 央企总部 10% 配额"],
        ])
    add_footer(s)

    # P9 招商漏斗
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P9 招商策略：四级漏斗 + 共享配套")
    add_body(s, [
        "L1 牌照：头部央企 / 行业协会 → 10% 配额（首批 3–5 家）",
        "L2 资本：追觅基金返投 + AI 腾讯生态导流 → 中型企业 20%",
        "L3 爬楼：仲量联行大数据触达 → 小型潮玩企业 40%",
        "L4 活动：科技开放麦 / 潮玩大赛带流 → 中小型服务机构 20%",
        "共享配套 10%（直播 / 选品 / 共享设计中心）→ 合资公司直营",
    ])
    add_footer(s)

    # P9b 市场租金对标
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P9b 市场租金对标：大零号湾 / 紫竹高新区 (2024-2025)")
    add_table(s, 1.2,
        ["园区 / 项目", "日租金 (元/㎡/天)", "类别"],
        [
            ["零号湾全球创新创业集聚区", "2.0–2.5", "标杆"],
            ["大零号湾科创成果转化中心", "2.0–2.5", "标杆"],
            ["华谊万创新所 / 上海人工智能产业园", "2.2", "主流"],
            ["紫竹信息数码港 (5A 甲级)", "2.1–2.5", "主流"],
            ["紫竹数字创意港 / 龙湖蓝海引擎", "1.5–3.0", "主流"],
            ["云境 443 / 夏日汇国际中心", "2.3–4.5", "高端"],
        ])
    add_body(s, [
        "结论:办公/产业研发用地主流区间 1.8-2.5 元/㎡/天 (年租 657-913 元/㎡)",
        "原 v1.0 测算 4.5-5.5 元偏高约一倍, v1.1 已下调至市场水平",
        "商业 1F 临街铺位可单独按 3.0-5.0 元定价, 不影响整体加权均值",
    ], top=4.5, height=2.0, font_size=14)
    add_footer(s)

    # P9c 科技企业服务中心
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P9c 科技企业服务中心：合资公司对外增值收入引擎")
    add_body(s, [
        "【定位】不是面向森马的费用,而是面向元谷入驻企业 + 大零号湾区域内潜在企业的 9 大类增值服务",
        "【选址】4# 楼 5F+ 或 5# 楼 5F+ 潮玩产业集群 (与北欧会客厅 / AI 共享设计中心同层联动)",
        "【9 大类服务】",
        "  ① 注册落户 ② 财税法 ③ 知识产权 ④ 政府补贴申报",
        "  ⑤ 人才与签证 ⑥ 投融资 (FA/并购) ⑦ 品牌与公关 (含出海)",
        "  ⑧ 数字化工具 (SaaS/AI 设计工作站) ⑨ 培训与认证",
        "【三年营收预测 (基础场景)】",
        "  T+1 年: 30 户 × 5 万 ≈ 150 万营收 / 净利 ≈ 53 万",
        "  T+2 年: 50 户 × 6 万 ≈ 300 万营收 / 净利 ≈ 114 万",
        "  T+3 年: 80 户 × 8 万 ≈ 640 万营收 / 净利 ≈ 256 万",
        "  三年累计营收 ≈ 1,090 万 / 累计净利 ≈ 423 万 / 丙方 30% 分红 ≈ 127 万",
    ], font_size=15)
    add_footer(s)

    # P10 楼宇地图
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P10 楼宇产业地图（1#–6#）")
    add_table(s, 1.2,
        ["楼栋", "既定功能", "我们的产业增益"],
        [
            ["1#", "森马总部 + 1–4F 零售", "福布斯榜单 / 国际首发节点"],
            ["2#", "二次元 Livehouse / 秀场", "科技开放麦 + 潮玩大赛主舞台"],
            ["3#", "1–4F 运动+萌宠 / 5F+ 酒店", "国际客户接待动线"],
            ["4#", "1–3F 潮玩艺术中心 / 4F 直播 / 5F+ 潮玩集群", "北欧会客厅 + AI 共享设计中心"],
            ["5#", "1–4F 动漫书店 / 5F+ 潮玩集群", "全国潮玩设计大赛主展区"],
            ["6#", "餐饮 / 商务宴请", "招商签约中心 + 央企 VIP"],
        ])
    add_footer(s)

    # P11 业态产能
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P11 重点业态产能配置")
    add_body(s, [
        "IP 潮玩选品 & 仓储式零售中心 ≈ 5,000㎡ → 合资公司直营，分润",
        "动漫潮玩谷主题街区 ≈ 3,000㎡ → 招商 + 联合运营",
        "潮玩艺术中心 ≈ 2,000㎡ → 北欧会客厅核心舱",
        "动漫主题书店 ≈ 1,500㎡ → IP 内容 + 阅读 + 活动",
        "森马展厅 & 二次元 Livehouse ≈ 700㎡ → 与品牌部联运",
    ])
    add_footer(s)

    # P12 合资公司
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P12 合作模式：合资公司架构")
    add_body(s, [
        "三方：森马（资产方） / 危总团队（属地资源 + LP） / 胡教授团队（运营 + 资源 GP）",
        "架构：合资公司持有元谷招商运营独家授权 + 共享配套直营权",
        "股比建议：森马 51% / 危总 19% / 胡教授团队 30%（含期权池）",
        "治理：3 人董事会，胡教授任首席战略运营官 (CSO)",
    ])
    add_footer(s)

    # P13 收入结构
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P13 收入结构 (v1.1 已含科技企业服务中心)")
    add_table(s, 1.2,
        ["收入类别", "测算逻辑", "支付方", "结算节奏"],
        [
            ["固定月费 Retainer", "20-35 万元/月 (推荐 28 万)", "甲方 (森马)", "按月"],
            ["招商佣金", "成交年租金的 1 / 1.5 / 2.5 个月", "甲方", "起租后"],
            ["活动运营收入", "单场净利分成", "外部+甲方", "按场"],
            ["专项奖项激励", "牌照 / 福布斯 / 小镇奖项一次性", "甲方", "按事件"],
            ["科技企业服务中心", "9 大类增值服务费", "外部入驻企业", "按服务"],
            ["股权分红", "合资公司净利 × 30% (含服务中心利润)", "合资公司", "年度"],
        ])
    add_footer(s)

    # P14 12 月里程碑
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P14 第一年关键里程碑（T+0 → T+12）")
    add_body(s, [
        "T+0 ~ T+3：合资公司注册 / 月费起算 / 北欧会客厅签约入驻 4# 楼",
        "T+3 ~ T+6：仲量联行数据上架 / 第一批爬楼名单 200 家 / 首场科技开放麦",
        "T+6 ~ T+9：两大产业牌照公开挂牌 / 追觅基金完成 1 期返投落地 5 家潮玩企业",
        "T+9 ~ T+12：全国潮玩设计大赛决赛 / 福布斯榜单年度发布 / 央企总部签约 1–2 家",
    ])
    add_footer(s)

    # P15 三年目标
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P15 三年产业目标")
    add_body(s, [
        "入驻潮玩相关企业 ≥120 家，匹配 10/10/20/20/40 五档配比",
        "合资公司直营业态合计 ≥10,000㎡，三年内现金流自给",
        "元谷品牌曝光：年度活动 ≥30 场、媒体声量 ≥5 亿次、政府奖项 ≥6 项",
        "推动元谷成为上海“科技时尚特色小镇”运营标杆",
    ])
    add_footer(s)

    # P16 风险
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P16 风险与对冲")
    add_table(s, 1.3,
        ["风险", "对冲机制"],
        [
            ["招商节奏不及预期", "月费保底 + 阶梯佣金双轨"],
            ["基金返投落地不顺", "AI 腾讯生态导流作为 Plan B"],
            ["国际外事波动", "北欧 + 日韩 + 东南亚多线布局"],
            ["政府政策调整", "央企 + 行业协会牌照前置锁定"],
        ])
    add_footer(s)

    # P17 团队
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P17 团队配备（招聘计划）")
    add_body(s, [
        "产业招商经理 × 1：仲量联行数据 + 爬楼 + 转化",
        "国际合作 & 活动策划 × 1：北欧会客厅 + 科技开放麦",
        "基金投后 & 政府关系 × 1：追觅基金返投 + 牌照申报",
        "胡教授（CSO，每周不少于 2 个工作日）",
        "1 名行政 / 财务（合资公司共享）",
    ])
    add_footer(s)

    # P18 我们的承诺
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P18 我们的承诺（写进协议）")
    add_body(s, [
        "北欧创新国际会客厅元谷站 12 个月内挂牌运营",
        "仲量联行爬楼数据 90 天内接入合资公司中台",
        "追觅基金 12 个月内完成首期返投落地 ≥3 家潮玩企业",
        "福布斯榜单 / 行业牌照 12 个月内挂牌 ≥2 项",
        "科技开放麦 12 个月内办 ≥10 场",
    ])
    add_footer(s)

    # P19 森马承诺
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P19 我们希望森马的承诺")
    add_body(s, [
        "授予合资公司独家招商运营权（5 年）",
        "商业可租面积 5.2 万㎡ 招商主导权",
        "共享配套（直播 / 选品 / 设计中心）由合资公司直营",
        "月费保底 + 招商佣金 + 奖项激励 + 股权分红组合",
        "政府关系联合体由森马牵头、合资公司执行",
    ])
    add_footer(s)

    # P20 测算速览
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P20 测算速览（详见 Excel）")
    add_table(s, 1.2,
        ["场景", "月费 (年化)", "招商佣金", "活动+奖项", "分红 (含服务中心)", "首年合计 (胡教授团队)"],
        [
            ["保守", "240 万", "72 万", "240 万", "72 万", "≈ 624 万元"],
            ["基础", "336 万", "132 万", "390 万", "144 万", "≈ 1,002 万元"],
            ["乐观", "420 万", "264 万", "560 万", "288 万", "≈ 1,532 万元"],
        ])
    add_body(s, [
        "另:合资公司层面新增「科技企业服务中心」首年营收 150 / 300 / 640 万元 (Sheet 05b)",
        "  其 38% 净利已通过分红 30% 反映在上表第 5 列, 不重复计入合计",
        "v1.1 已按真实市场租金 (1.8-2.5 元/㎡/天) 矫正; 详见 02_测算模型/合作收益测算模型.xlsx",
    ], top=4.4, height=2.4, font_size=13)
    add_footer(s)

    # P21 下一步
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "P21 下一步行动")
    add_body(s, [
        "本周：双方就月费区间 + 股比锁定意向",
        "次周：签订 MOU（基于本协议草案）",
        "30 天内：合资公司注册 + 三人核心团队就位",
        "60 天内：北欧会客厅元谷站 + 第一场科技开放麦同步启动",
    ])
    add_footer(s)

    # P22 收尾
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = PRIMARY; bg.line.fill.background()
    box = s.shapes.add_textbox(Inches(1.0), Inches(2.6), Inches(11.3), Inches(3.0)).text_frame
    box.word_wrap = True
    p1 = box.paragraphs[0]; p1.text = "让元谷成为中国潮玩产业"; p1.font.size = Pt(54); p1.font.bold = True; p1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p2 = box.add_paragraph(); p2.text = "走向世界的国际客厅。"; p2.font.size = Pt(54); p2.font.bold = True; p2.font.color.rgb = ACCENT
    p3 = box.add_paragraph(); p3.text = ""
    p4 = box.add_paragraph(); p4.text = "胡教授团队 × 森马集团 联合运营合资公司"; p4.font.size = Pt(20); p4.font.color.rgb = RGBColor(0xCF, 0xD8, 0xE3)

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
