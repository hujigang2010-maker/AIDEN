"""会议纪要后 · 报价 PPT + 变更 PPT + 变更 Excel

依据《产业园区合作方案讨论会》纪要/录音整理:
  - 招商佣金: 统一 1.5 个月租金(渠道成交)
  - 沙龙: 先试点 3 场一期, 设落地考核, 再决定续做
  - 媒体: 可选; 打包示意 12 万 / 传播量 ≥300 万
  - 挂牌: 至少挂 1 项(科协/科委相关), 约 10 万/项
  - 超额奖励: 不按整盘 90% 粗放计提; 按我方贡献细化
  - 政策服务: 正式单列(园区专项一次性激励 + 企业兑现抽成)

产出:
  1) 元谷合作服务报价_会议纪要版.pptx   ← 主交付(给对方看的报价)
  2) 元谷合作条款_会议调整内容.pptx     ← 仅含因纪要而调整的部分
  3) 元谷合作条款_会议调整内容.xlsx     ← 同上 Excel
"""
from __future__ import annotations
from pathlib import Path

from openpyxl import Workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

DIR = Path(__file__).resolve().parent
OUT_QUOTE = DIR / "元谷合作服务报价_会议纪要版.pptx"
OUT_CHG_PPT = DIR / "元谷合作条款_会议调整内容.pptx"
OUT_CHG_XLS = DIR / "元谷合作条款_会议调整内容.xlsx"

PRIMARY = RGBColor(0x0F, 0x24, 0x4E)
ACCENT = RGBColor(0xF2, 0x7E, 0x2D)
GOLD = RGBColor(0xC9, 0xA2, 0x4B)
TEAL = RGBColor(0x18, 0x8B, 0x8B)
GREEN = RGBColor(0x1E, 0x8E, 0x5A)
LIGHT = RGBColor(0xF4, 0xF6, 0xFA)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD8, 0xDF, 0xEA)
TEXT = RGBColor(0x21, 0x2B, 0x42)
MUTED = RGBColor(0x66, 0x70, 0x86)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

_page = {"n": 0}


def _font(run, size, *, bold=False, color=TEXT, name="微软雅黑"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    from pptx.oxml.ns import qn
    rpr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea"):
        el = rpr.find(qn(tag))
        if el is None:
            el = rpr.makeelement(qn(tag), {})
            rpr.append(el)
        el.set("typeface", name)


def _rect(slide, l, t, w, h, fill=None, line=None, line_w=1.0):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def _text(slide, l, t, w, h, lines, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    for i, spec in enumerate(lines):
        text, size, bold, color = spec[0], spec[1], spec[2], spec[3]
        sa = spec[4] if len(spec) > 4 else 3
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(sa)
        r = p.add_run()
        r.text = text
        _font(r, size, bold=bold, color=color)
    return tb


def bg(slide):
    _rect(slide, -0.1, -0.1, 13.6, 7.7, fill=LIGHT)


def header(slide, kicker, title, right="报价 · 会议纪要版"):
    _rect(slide, 0, 0, 13.333, 1.15, fill=PRIMARY)
    _rect(slide, 0, 1.15, 13.333, 0.06, fill=ACCENT)
    _rect(slide, 0.55, 0.26, 0.12, 0.62, fill=ACCENT)
    _text(slide, 0.85, 0.20, 9.5, 0.34, [(kicker, 11, False, RGBColor(0xB9, 0xC6, 0xDC))])
    _text(slide, 0.85, 0.46, 9.8, 0.55, [(title, 20, True, WHITE)])
    _text(
        slide,
        10.0,
        0.22,
        2.9,
        0.75,
        [("元谷 YUANGU", 12, True, RGBColor(0xE8, 0xC8, 0x7A)), (right, 9, False, RGBColor(0xB9, 0xC6, 0xDC))],
        align=PP_ALIGN.RIGHT,
    )


def footer(slide, label="元谷合作服务报价 · 会议纪要版 · 仅供商务沟通"):
    _page["n"] += 1
    _rect(slide, 0, 7.28, 13.333, 0.22, fill=PRIMARY)
    _text(slide, 0.55, 7.28, 10.2, 0.22, [(label, 8, False, RGBColor(0xC6, 0xD0, 0xE2))], anchor=MSO_ANCHOR.MIDDLE)
    _text(
        slide,
        11.7,
        7.28,
        1.1,
        0.22,
        [(f"{_page['n']:02d}", 8, True, RGBColor(0xE8, 0xC8, 0x7A))],
        align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def chip(slide, l, t, w, h, text, *, fill=ACCENT, size=12):
    _rect(slide, l, t, w, h, fill=fill)
    _text(slide, l, t, w, h, [(text, size, True, WHITE)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def kpi(slide, l, t, w, h, label, value, sub, *, accent=ACCENT):
    _rect(slide, l, t, w, h, fill=CARD, line=LINE, line_w=0.75)
    _rect(slide, l, t, 0.1, h, fill=accent)
    _text(slide, l + 0.22, t + 0.14, w - 0.35, 0.28, [(label, 10.5, False, MUTED)])
    _text(slide, l + 0.22, t + 0.42, w - 0.35, 0.5, [(value, 20, True, PRIMARY)])
    _text(slide, l + 0.22, t + 1.0, w - 0.35, 0.4, [(sub, 11, False, MUTED)])


def bullets(slide, l, t, w, h, items, *, fs=13, gap=8, color=TEXT, marker="●", mc=ACCENT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r0 = p.add_run()
        r0.text = f"{marker} "
        _font(r0, fs, bold=True, color=mc)
        r1 = p.add_run()
        r1.text = item
        _font(r1, fs, bold=False, color=color)


def table(slide, l, t, w, headers, rows, *, col_w=None, row_h=0.55, fs=11, accent_first=False):
    n_cols = len(headers)
    n_rows = 1 + len(rows)
    if col_w is None:
        col_w = [w / n_cols] * n_cols
    shape = slide.shapes.add_table(n_rows, n_cols, Inches(l), Inches(t), Inches(w), Inches(row_h * n_rows))
    tbl = shape.table
    for i, cw in enumerate(col_w):
        tbl.columns[i].width = Inches(cw)
    for j, htxt in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = htxt
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                _font(r, fs, bold=True, color=WHITE)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if j == 0 else PP_ALIGN.LEFT
                for r in p.runs:
                    bold = accent_first and j == 0
                    _font(r, fs, bold=bold, color=PRIMARY if bold else TEXT)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xEF, 0xF2, 0xF8) if i % 2 else CARD
    return shape


def build_quote_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    _page["n"] = 0

    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, 13.333, 7.5, fill=PRIMARY)
    _rect(s, 0, 0, 0.28, 7.5, fill=ACCENT)
    chip(s, 1.1, 1.3, 3.2, 0.42, "会议纪要后报价版", fill=ACCENT, size=12)
    _text(s, 1.1, 2.0, 11, 0.4, [("森马（上海）国际运营中心 · 元谷项目", 14, False, RGBColor(0xB9, 0xC6, 0xDC))])
    _text(s, 1.1, 2.5, 11, 1.0, [("合作服务报价", 36, True, WHITE)])
    _text(
        s,
        1.1,
        3.6,
        11,
        0.8,
        [
            ("轻启动 · 阶段验证 · 结果导向", 18, False, RGBColor(0xE8, 0xC8, 0x7A)),
            ("沙龙试点 / 招商佣金 / 挂牌 / 政策服务 / 媒体可选", 14, False, RGBColor(0xCF, 0xD8, 0xE8)),
        ],
    )
    _text(
        s,
        1.1,
        5.3,
        11,
        1.0,
        [
            ("胡教授团队｜复旦大学住房政策研究中心 / 上海市科技企业联合会相关资源", 13, False, RGBColor(0xCF, 0xD8, 0xE8)),
            ("呈：森马资管相关负责人  ·  第三方专业服务报价", 12, False, RGBColor(0x9F, 0xB2, 0xD0)),
        ],
    )
    _text(s, 1.1, 6.7, 11, 0.4, [("依据合作方案讨论会沟通共识整理 · 仅供商务沟通", 11, False, RGBColor(0x7A, 0x8C, 0xA8))])

    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "01 · 报价总览", "本期合作：先试点、再放大")
    kpi(s, 0.7, 1.4, 3.0, 1.55, "沙龙试点", "3 场 / 一期", "建议执行费 15 万元", accent=ACCENT)
    kpi(s, 3.9, 1.4, 3.0, 1.55, "招商佣金", "1.5 个月", "渠道成交年租金", accent=TEAL)
    kpi(s, 7.1, 1.4, 2.9, 1.55, "挂牌(建议)", "至少 1 项", "约 10 万元/项", accent=GOLD)
    kpi(s, 10.2, 1.4, 2.45, 1.55, "启动现金", "约 25 万", "3 场沙龙+1 项挂牌", accent=GREEN)
    table(
        s,
        0.7,
        3.2,
        11.95,
        ["服务模块", "本期报价 / 标准", "必选/可选", "结算节点"],
        [
            ["沙龙执行（试点一期）", "3 场打包 15 万元（可按场结算）", "建议必选", "启动付 50% / 第 3 场后付 50%"],
            ["招商佣金", "实际成交年租金 × 1.5 个月", "成交后结算", "起租后 30 日内"],
            ["挂牌服务", "10 万元/项；本期建议至少 1 项", "建议至少 1 项", "挂牌前一次性"],
            ["政策服务", "园区专项落地奖 5–10 万；企业兑现按执行额抽成", "建议纳入", "成果确认后结算"],
            ["媒体支持", "打包示意 12 万元（≥300 万传播量）", "可选", "按选定包结算"],
            ["超额奖励", "按我方贡献成交面积细化，不按整盘粗放计提", "可选/后议", "阶段考核后"],
        ],
        col_w=[2.8, 4.5, 1.9, 2.75],
        row_h=0.52,
        fs=11,
        accent_first=True,
    )
    footer(s)

    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "02 · 合作原则", "费用按“招商基础设施投入”理解，节奏按“阶段验证”推进")
    cards = [
        ("轻启动", "不一口气买断全年", "先做 3 场沙龙试点，再决定是否续做后三场。", ACCENT),
        ("结果导向", "设明确落地考核", "试点期考核：至少 1 家企业落地，再推进下一期。", TEAL),
        ("佣金平衡", "统一 1.5 个月", "符合市场水平，兼顾积极性与租金竞争力。", GOLD),
        ("政策单列", "不只是附赠", "园区专项政策 + 企业政策兑现，设成果激励。", GREEN),
    ]
    for i, (a, b, c, ac) in enumerate(cards):
        x = 0.7 + i * 3.15
        _rect(s, x, 1.5, 3.0, 4.5, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, x, 1.5, 3.0, 1.1, fill=ac)
        _text(s, x, 1.6, 3.0, 0.4, [(a, 12, False, RGBColor(0xFF, 0xF0, 0xE0))], align=PP_ALIGN.CENTER)
        _text(s, x + 0.15, 2.05, 2.7, 0.45, [(b, 15, True, WHITE)], align=PP_ALIGN.CENTER)
        _text(s, x + 0.2, 2.95, 2.6, 2.6, [(c, 14, False, TEXT)])
    footer(s)

    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "03 · 沙龙执行报价", "试点一期：3 场")
    kpi(s, 0.7, 1.4, 3.85, 1.5, "试点场次", "3 场", "不以 6 场一次性打包启动", accent=ACCENT)
    kpi(s, 4.75, 1.4, 3.85, 1.5, "建议执行费", "15 万元", "按原 6 场 30 万对半折算", accent=TEAL)
    kpi(s, 8.8, 1.4, 3.85, 1.5, "试点考核", "≥1 家落地", "完成后再启动下一期", accent=GREEN)
    table(
        s,
        0.7,
        3.15,
        11.95,
        ["项目", "内容", "说明"],
        [
            ["单场目标", "到场目标产业客户 ≥30 家", "企业家/机构/协会/校友等高信任客群"],
            ["费用构成", "沙龙策划、邀约、主持、转化跟进", "场地与基础会务保障由园区侧配合"],
            ["支付方式", "启动付 50%（7.5 万）/ 第 3 场完成后付 50%", "也可按场 5 万元结算"],
            ["续做条件", "试点期至少 1 家企业落地", "达标后启动后 3 场（费用另议/同比）"],
            ["历史参考", "今年约 3 场 · 邀约约 120 人 · 成交约 3 户 · 约 2,000㎡", "用于理解转化节奏，不作为承诺"],
        ],
        col_w=[2.2, 5.0, 4.75],
        row_h=0.55,
        accent_first=True,
    )
    footer(s)

    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "04 · 招商佣金报价", "统一按 1.5 个月租金结算")
    _rect(s, 0.7, 1.4, 11.95, 1.7, fill=CARD, line=LINE, line_w=0.75)
    _rect(s, 0.7, 1.4, 0.12, 1.7, fill=ACCENT)
    _text(s, 1.05, 1.55, 11.2, 0.4, [("结算标准", 14, True, PRIMARY)])
    _text(
        s,
        1.05,
        2.1,
        11.2,
        0.8,
        [
            ("招商佣金 = 实际成交年租金 × 1.5 个月", 20, True, ACCENT),
            ("年租金 = 成交单价 × 面积 × 365（以合同载明为准）", 13, False, MUTED),
        ],
    )
    table(
        s,
        0.7,
        3.35,
        11.95,
        ["规则", "口径"],
        [
            ["适用对象", "经我方渠道提报、跟进并促成成交的客户"],
            ["不适用", "园区自有团队、已备案的既有意向客户、政府直接导入且我方未实质跟进的客户"],
            ["支付时点", "租赁合同签订且租户起租后 30 日内"],
            ["租金平衡", "佣金纳入整体租金成本考量；统一 1.5 个月，避免推高租金影响去化"],
            ["测算示例", "若单户年租金 200 万，则佣金约 25 万；10 户量级即可显著覆盖试点投入"],
        ],
        col_w=[2.4, 9.55],
        row_h=0.55,
        accent_first=True,
    )
    footer(s)

    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "05 · 挂牌与媒体报价", "挂牌建议启动；媒体按需选配")
    _rect(s, 0.7, 1.4, 6.0, 5.4, fill=CARD, line=LINE, line_w=0.75)
    _rect(s, 0.7, 1.4, 6.0, 0.55, fill=PRIMARY)
    _text(s, 0.95, 1.4, 5.5, 0.55, [("挂牌服务", 15, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
    bullets(
        s,
        0.95,
        2.2,
        5.5,
        4.2,
        [
            "标准：10 万元/项",
            "本期建议：至少挂 1 项",
            "优先方向：上海市科协/科委相关牌照（如科技企业服务中心）",
            "价值：提升园区信誉度与企业服务能力",
            "支付：挂牌前一次性付清",
            "有效性：长期挂牌使用（按牌照规则执行）",
        ],
        fs=13.5,
        gap=10,
    )
    _rect(s, 6.95, 1.4, 5.7, 5.4, fill=CARD, line=LINE, line_w=0.75)
    _rect(s, 6.95, 1.4, 5.7, 0.55, fill=TEAL)
    _text(s, 7.2, 1.4, 5.2, 0.55, [("媒体支持（可选）", 15, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
    bullets(
        s,
        7.2,
        2.2,
        5.2,
        4.2,
        [
            "打包示意：12 万元",
            "内容：中央级/全国级发稿约 10 篇 + 原创内容 3 天",
            "传播量承诺：不低于 300 万",
            "定位：可选项；园区自有媒体团队可自行统筹",
            "选用原则：符合园区整体推广策略",
            "评估依据：是否服务实际招商转化",
        ],
        fs=13.5,
        gap=10,
    )
    footer(s)

    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "06 · 政策服务报价", "正式单列：园区专项 + 企业兑现")
    table(
        s,
        0.7,
        1.4,
        11.95,
        ["政策服务类型", "服务内容", "费用 / 激励"],
        [
            ["园区专项产业政策", "协助园区与政府沟通、推动专项政策落地", "成功落地后一次性奖励 5–10 万元（按支持力度协商）"],
            ["入园企业政策兑现", "协助企业申报、落地与兑现（补贴/资质等）", "按最终执行金额抽成（比例另议）"],
            ["日常政策陪跑", "在沙龙/招商服务过程中的政策识别与路径辅导", "纳入试点合作基础服务"],
        ],
        col_w=[3.0, 4.4, 4.55],
        row_h=0.95,
        accent_first=True,
    )
    bullets(
        s,
        0.7,
        4.7,
        12,
        2.0,
        [
            "政策有时效性（常见 3–5 年周期），需在窗口期内推进。",
            "企业侧常有门槛条件（如人工成本占比等），需政企双方共同努力。",
            "不做“包拿补贴”承诺；以合规申报、路径清晰、成果可核验为准。",
        ],
        fs=13.5,
        gap=8,
    )
    footer(s)

    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "07 · 建议启动组合", "先付得起、验得清、推得动")
    table(
        s,
        0.7,
        1.4,
        11.95,
        ["启动项", "金额", "说明"],
        [
            ["沙龙试点 3 场", "15 万元", "启动付 7.5 万 + 第 3 场后付 7.5 万"],
            ["挂牌 1 项（建议）", "10 万元", "优先科协/科委相关牌照"],
            ["启动现金小计", "25 万元", "不含佣金、政策奖、媒体"],
            ["招商佣金", "成交后按 1.5 个月", "不占用启动现金；成功后结算"],
            ["政策专项落地奖", "5–10 万元/次", "有成果再付"],
            ["媒体支持", "12 万元（可选）", "按园区市场部需求决定是否采购"],
        ],
        col_w=[3.2, 3.0, 5.75],
        row_h=0.58,
        accent_first=True,
    )
    chip(
        s,
        0.7,
        5.5,
        11.95,
        0.7,
        "建议口径：森马无需一次性买断全年；先用约 25 万启动试点，用首批落地验证合作价值。",
        fill=PRIMARY,
        size=14,
    )
    _text(
        s,
        0.7,
        6.4,
        12,
        0.4,
        [("超额奖励不按整盘出租率 90% 粗放计提；若设置，按我方贡献成交面积/转化结果细化后议。", 12.5, True, ACCENT)],
    )
    footer(s)

    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "08 · 阶段节奏与考核", "三步推进")
    steps = [
        ("第 1 步", "确认合作口径", "本周", "定位、授权边界、宣传口径、政策服务边界、客户归属规则。", ACCENT),
        ("第 2 步", "启动试点一期", "2–3 周内", "启动 3 场沙龙 +（建议）1 项挂牌；形成首批名单与复访计划。", TEAL),
        ("第 3 步", "复访与续做决策", "30–90 天", "至少 1 家落地后启动下一期；同步推进佣金结算与政策成果。", GREEN),
    ]
    for i, (no, title, when, body, ac) in enumerate(steps):
        x = 0.7 + i * 4.15
        _rect(s, x, 1.45, 3.95, 4.3, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, x, 1.45, 3.95, 1.2, fill=ac)
        _text(s, x + 0.2, 1.55, 3.55, 0.35, [(no, 12, False, RGBColor(0xFF, 0xF0, 0xE0))], align=PP_ALIGN.CENTER)
        _text(s, x + 0.2, 1.95, 3.55, 0.5, [(title, 18, True, WHITE)], align=PP_ALIGN.CENTER)
        chip(s, x + 0.9, 2.95, 2.1, 0.4, when, fill=PRIMARY, size=12)
        _text(s, x + 0.25, 3.6, 3.45, 1.8, [(body, 14, False, TEXT)])
    footer(s)

    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, 13.333, 7.5, fill=PRIMARY)
    _rect(s, 0, 0, 0.28, 7.5, fill=ACCENT)
    _text(s, 1.1, 2.3, 11, 0.4, [("报价结论", 14, False, RGBColor(0xE8, 0xC8, 0x7A))])
    _text(s, 1.1, 2.9, 11, 1.2, [("先试点 3 场 + 挂牌 1 项", 28, True, WHITE), ("佣金 1.5 个月 · 政策成果另计", 28, True, WHITE)])
    bullets(
        s,
        1.1,
        4.6,
        11,
        1.8,
        [
            "启动现金约 25 万元（沙龙 15 万 + 挂牌 10 万）；",
            "媒体 12 万可选；政策专项落地奖 5–10 万/次；企业兑现抽成另议；",
            "试点达标后再放大，控制风险、匹配产出。",
        ],
        fs=15,
        gap=10,
        color=RGBColor(0xE5, 0xEB, 0xF5),
        marker="▸",
        mc=RGBColor(0xE8, 0xC8, 0x7A),
    )

    prs.save(OUT_QUOTE)
    print(f"Wrote {OUT_QUOTE} pages={_page['n']}+cover/end")


def build_change_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    _page["n"] = 0

    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, 13.333, 7.5, fill=PRIMARY)
    _rect(s, 0, 0, 0.28, 7.5, fill=ACCENT)
    _text(s, 1.1, 2.2, 11, 0.4, [("元谷合作条款", 14, False, RGBColor(0xE8, 0xC8, 0x7A))])
    _text(s, 1.1, 2.8, 11, 1.0, [("会议沟通后调整内容", 32, True, WHITE)])
    _text(
        s,
        1.1,
        4.1,
        11,
        1.0,
        [
            ("仅整理本次会议形成的条款调整，便于单独审阅。", 15, False, RGBColor(0xCF, 0xD8, 0xE8)),
            ("是否并入整体方案 PPT，由贵方决定。", 14, False, RGBColor(0x9F, 0xB2, 0xD0)),
        ],
    )

    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "调整内容一览", "六项条款", right="会议调整内容")
    table(
        s,
        0.7,
        1.4,
        11.95,
        ["条款", "调整后口径"],
        [
            ["招商佣金", "统一按实际成交年租金的 1.5 个月；仅适用我方渠道成交客户"],
            ["沙龙执行", "先做 3 场试点一期；建议执行费 15 万元；试点考核至少 1 家落地"],
            ["媒体支持", "可选；打包示意 12 万元（约 10 篇 + 原创 3 天，传播量 ≥300 万）"],
            ["挂牌服务", "10 万元/项；本期建议至少挂 1 项（优先科协/科委相关）"],
            ["超额奖励", "不按整盘出租率 90% 粗放计提；按我方贡献成交结果细化后议"],
            ["政策服务", "正式单列：园区专项落地奖 5–10 万；企业政策兑现按执行额抽成"],
        ],
        col_w=[2.6, 9.35],
        row_h=0.72,
        accent_first=True,
    )
    footer(s, "元谷合作条款 · 会议调整内容 · 仅供商务沟通")

    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "沙龙执行", "试点一期 3 场 / 建议 15 万元", right="会议调整内容")
    kpi(s, 0.7, 1.4, 3.85, 1.55, "场次", "3 场", "一期试点", accent=ACCENT)
    kpi(s, 4.75, 1.4, 3.85, 1.55, "执行费", "15 万元", "可 50%/50% 支付", accent=TEAL)
    kpi(s, 8.8, 1.4, 3.85, 1.55, "考核", "≥1 家落地", "达标后续做", accent=GREEN)
    bullets(
        s,
        0.7,
        3.3,
        12,
        3.3,
        [
            "不以 6 场、30 万元一次性打包作为启动条件。",
            "单场目标产业客户 ≥30 家；场地与基础会务由园区配合。",
            "历史参考：约 3 场、邀约约 120 人、成交约 3 户、约 2,000㎡。",
            "试点完成后，根据转化效果决定是否启动下一期。",
        ],
        fs=14,
        gap=12,
    )
    footer(s, "元谷合作条款 · 会议调整内容 · 仅供商务沟通")

    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "招商佣金", "统一 1.5 个月租金", right="会议调整内容")
    _rect(s, 0.7, 1.45, 11.95, 1.5, fill=CARD, line=LINE, line_w=0.75)
    _text(s, 1.0, 1.7, 11.4, 0.4, [("招商佣金 = 实际成交年租金 × 1.5 个月", 22, True, ACCENT)])
    _text(s, 1.0, 2.3, 11.4, 0.4, [("取消按面积分档的 1.5 / 1.75 / 2.0 个月结构", 14, False, MUTED)])
    bullets(
        s,
        0.7,
        3.3,
        12,
        3.3,
        [
            "仅对我方渠道提报并促成成交的客户计提。",
            "园区自招、既有意向、政府直导且我方未实质跟进的客户，不计提。",
            "佣金需纳入整体租金成本平衡，避免推高租金、增加招商难度。",
            "起租后 30 日内结算。",
        ],
        fs=14,
        gap=12,
    )
    footer(s, "元谷合作条款 · 会议调整内容 · 仅供商务沟通")

    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "媒体 / 挂牌 / 超额奖励", "三项口径", right="会议调整内容")
    blocks = [
        ("媒体支持", "可选模块", "打包示意 12 万元：约 10 篇中央级/全国级发稿 + 原创内容 3 天；传播量不低于 300 万。园区自有媒体团队可自行统筹，是否采购以招商实效评估。", TEAL),
        ("挂牌服务", "建议至少 1 项", "10 万元/项。优先上海市科协/科委相关牌照（如科技企业服务中心）。挂牌前一次性支付，长期使用按牌照规则。", ACCENT),
        ("超额奖励", "细化后议", "不采用“整盘出租率 ≥90% 即奖”的粗放方式（非独家情形下不合理）。若设置，按我方贡献成交面积或转化结果另行约定。", GOLD),
    ]
    for i, (a, b, c, ac) in enumerate(blocks):
        y = 1.4 + i * 1.8
        _rect(s, 0.7, y, 11.95, 1.65, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, 0.7, y, 2.4, 1.65, fill=ac)
        _text(s, 0.7, y + 0.35, 2.4, 0.4, [(a, 14, True, WHITE)], align=PP_ALIGN.CENTER)
        _text(s, 0.7, y + 0.85, 2.4, 0.4, [(b, 12, False, RGBColor(0xFF, 0xF0, 0xE0))], align=PP_ALIGN.CENTER)
        _text(s, 3.4, y + 0.35, 8.9, 1.1, [(c, 13.5, False, TEXT)], anchor=MSO_ANCHOR.MIDDLE)
    footer(s, "元谷合作条款 · 会议调整内容 · 仅供商务沟通")

    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "政策服务", "正式纳入合作框架", right="会议调整内容")
    table(
        s,
        0.7,
        1.4,
        11.95,
        ["类型", "内容", "激励"],
        [
            ["园区专项产业政策", "协助园区争取/落地政府专项政策", "成功落地后一次性奖励 5–10 万元"],
            ["入园企业政策兑现", "协助企业申报与兑现，服务不止于签约", "按最终执行金额抽成（比例另议）"],
            ["基础政策陪跑", "招商与沙龙过程中的政策识别、路径辅导", "纳入试点基础服务"],
        ],
        col_w=[3.2, 5.0, 3.75],
        row_h=0.9,
        accent_first=True,
    )
    bullets(
        s,
        0.7,
        4.6,
        12,
        2.0,
        [
            "政策窗口常见 3–5 年周期，需在有效期内推进。",
            "企业侧存在申报门槛（如人工成本占比等），属双方共同努力事项。",
            "以合规申报与可核验成果为准，不做包拿承诺。",
        ],
        fs=13.5,
        gap=8,
    )
    footer(s, "元谷合作条款 · 会议调整内容 · 仅供商务沟通")

    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, "建议启动组合", "现金与后结算分开", right="会议调整内容")
    table(
        s,
        0.7,
        1.45,
        11.95,
        ["项目", "金额 / 标准", "时点"],
        [
            ["沙龙试点 3 场", "15 万元", "启动 50% + 第 3 场后 50%"],
            ["挂牌 1 项", "10 万元", "挂牌前一次性"],
            ["启动现金小计", "25 万元", "—"],
            ["招商佣金", "年租金 × 1.5 个月", "起租后 30 日内"],
            ["政策专项落地奖", "5–10 万元/次", "成果确认后"],
            ["企业政策兑现抽成", "按执行额比例（另议）", "兑现完成后"],
            ["媒体支持", "12 万元（可选）", "按选定包结算"],
        ],
        col_w=[3.4, 4.4, 4.15],
        row_h=0.55,
        accent_first=True,
    )
    footer(s, "元谷合作条款 · 会议调整内容 · 仅供商务沟通")

    prs.save(OUT_CHG_PPT)
    print(f"Wrote {OUT_CHG_PPT} pages={_page['n']}+cover")


def build_change_xlsx():
    wb = Workbook()
    thin = Side(style="thin", color="B0BEC5")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    fill_p = PatternFill("solid", fgColor="0F244E")
    fill_l = PatternFill("solid", fgColor="EAEEF5")
    fill_alt = PatternFill("solid", fgColor="F4F6FA")
    h2 = Font(name="微软雅黑", size=11, bold=True, color="FFFFFF")
    body = Font(name="微软雅黑", size=10, color="212B42")
    center = Alignment(horizontal="center", vertical="center", wrap_text=True)
    left = Alignment(horizontal="left", vertical="center", wrap_text=True)

    def widths(ws, ws_list):
        for i, w in enumerate(ws_list, 1):
            ws.column_dimensions[get_column_letter(i)].width = w

    def title(ws, text, span):
        ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=span)
        c = ws.cell(row=1, column=1, value=text)
        c.font = Font(name="微软雅黑", size=16, bold=True, color="FFFFFF")
        c.fill = fill_p
        c.alignment = Alignment(horizontal="left", vertical="center", indent=1)
        ws.row_dimensions[1].height = 34

    def hdr(ws, row, headers):
        for j, h in enumerate(headers, 1):
            c = ws.cell(row=row, column=j, value=h)
            c.font = h2
            c.fill = fill_p
            c.alignment = center
            c.border = border
        ws.row_dimensions[row].height = 26

    def put(ws, r, vals):
        for j, v in enumerate(vals, 1):
            c = ws.cell(row=r, column=j, value=v)
            c.font = body
            c.alignment = left
            c.border = border
            c.fill = fill_alt if r % 2 == 0 else fill_l

    ws = wb.active
    ws.title = "01 调整后条款"
    title(ws, "元谷合作条款 · 会议沟通后调整内容", 4)
    widths(ws, [6, 18, 55, 28])
    hdr(ws, 3, ["序号", "条款", "调整后口径", "备注"])
    rows = [
        (1, "招商佣金", "实际成交年租金 × 1.5 个月；仅适用我方渠道成交客户", "取消面积分档 1.5/1.75/2.0"),
        (2, "沙龙执行", "试点一期 3 场；建议执行费 15 万元；考核至少 1 家落地后再续做", "不以 6 场/30 万一次性启动"),
        (3, "媒体支持", "可选；打包示意 12 万元（约10篇+原创3天，传播量≥300万）", "园区自有媒体可统筹"),
        (4, "挂牌服务", "10 万元/项；本期建议至少 1 项（优先科协/科委相关）", "挂牌前一次性支付"),
        (5, "超额奖励", "不按整盘出租率90%粗放计提；按我方贡献成交结果细化后议", "非独家情形下调整"),
        (6, "政策服务", "正式单列：园区专项落地奖5–10万；企业兑现按执行额抽成", "由附赠改为可量化激励"),
    ]
    for i, r in enumerate(rows, start=4):
        put(ws, i, list(r))

    ws = wb.create_sheet("02 报价明细")
    title(ws, "报价明细（会议纪要版）", 5)
    widths(ws, [6, 22, 22, 18, 40])
    hdr(ws, 3, ["序号", "服务模块", "金额/标准", "必选/可选", "结算方式"])
    detail = [
        (1, "沙龙试点（3场）", "150,000 元", "建议必选", "启动50% + 第3场后50%；或5万/场"),
        (2, "挂牌（建议1项）", "100,000 元/项", "建议至少1项", "挂牌前一次性"),
        (3, "启动现金小计", "250,000 元", "—", "沙龙15万+挂牌10万"),
        (4, "招商佣金", "年租金×1.5个月", "成交后结算", "起租后30日内"),
        (5, "园区专项政策落地奖", "50,000–100,000 元/次", "建议纳入", "成果确认后一次性"),
        (6, "企业政策兑现抽成", "按执行金额比例（另议）", "建议纳入", "兑现完成后结算"),
        (7, "媒体支持", "120,000 元", "可选", "按选定服务包结算"),
        (8, "超额奖励", "按我方贡献细化（后议）", "可选/后议", "阶段考核后"),
    ]
    for i, r in enumerate(detail, start=4):
        put(ws, i, list(r))
    ws.cell(row=13, column=2, value="说明：沙龙场地与基础会务保障由园区侧配合；政策支持以申报条件、主管部门审核与项目合规为准。").font = Font(
        name="微软雅黑", size=9, italic=True, color="667086"
    )
    ws.merge_cells(start_row=13, start_column=2, end_row=13, end_column=5)

    ws = wb.create_sheet("03 沙龙试点")
    title(ws, "沙龙试点一期安排", 4)
    widths(ws, [8, 22, 40, 30])
    hdr(ws, 3, ["项目", "内容", "标准", "说明"])
    salon = [
        ("场次", "试点一期", "3 场", "先验证再放大"),
        ("执行费", "打包", "150,000 元", "可对半支付或按场5万"),
        ("单场客群", "目标产业客户", "≥30 家/场", "企业家/机构/协会/校友等"),
        ("考核", "落地转化", "至少 1 家企业落地", "达标后续做下一期"),
        ("历史参考", "今年数据", "约3场 / 邀约约120人 / 成交约3户 / 约2000㎡", "理解节奏，不作承诺"),
        ("园区配合", "会务保障", "场地、基础人员保障", "园区承担基础会务成本"),
    ]
    for i, r in enumerate(salon, start=4):
        put(ws, i, list(r))

    ws = wb.create_sheet("04 政策服务")
    title(ws, "政策服务条款", 4)
    widths(ws, [8, 24, 45, 30])
    hdr(ws, 3, ["序号", "类型", "服务内容", "激励方式"])
    policy = [
        (1, "园区专项产业政策", "协助园区与政府沟通，推动专项产业政策落地", "成功落地一次性奖励5–10万元"),
        (2, "入园企业政策兑现", "协助企业完成申报、落地与兑现", "按最终执行金额抽成（比例另议）"),
        (3, "基础政策陪跑", "招商/沙龙过程中的政策识别与路径辅导", "纳入试点基础服务"),
    ]
    for i, r in enumerate(policy, start=4):
        put(ws, i, list(r))
    ws.cell(row=8, column=2, value="备注：政策常见3–5年周期；企业申报常有门槛条件；不做包拿承诺。").font = Font(
        name="微软雅黑", size=9, italic=True, color="667086"
    )
    ws.merge_cells(start_row=8, start_column=2, end_row=8, end_column=4)

    ws = wb.create_sheet("05 启动组合")
    title(ws, "建议启动组合与现金测算", 4)
    widths(ws, [8, 28, 18, 40])
    hdr(ws, 3, ["序号", "项目", "金额（元）", "说明"])
    start = [
        (1, "沙龙试点3场", 150000, "建议必选"),
        (2, "挂牌1项", 100000, "建议至少1项"),
        (3, "启动现金合计", 250000, "不含佣金/政策奖/媒体"),
        (4, "媒体支持（可选加项）", 120000, "按需"),
        (5, "若含媒体的启动上限", 370000, "25万+12万"),
    ]
    for i, r in enumerate(start, start=4):
        put(ws, i, list(r))
        ws.cell(row=i, column=3).number_format = "#,##0"
        ws.cell(row=i, column=3).alignment = Alignment(horizontal="right", vertical="center")

    ws = wb.create_sheet("06 原口径对照")
    title(ws, "原方案口径 → 会议后口径（对照表）", 4)
    widths(ws, [16, 40, 40, 28])
    hdr(ws, 3, ["条款", "原方案口径", "会议后口径", "影响"])
    cmp_rows = [
        ("招商佣金", "1.5/1.75/2.0个月（按面积分档）", "统一1.5个月；仅渠道成交", "结构简化，利于租金平衡"),
        ("沙龙执行", "6场打包30万，建议一次性", "先3场试点/建议15万；≥1家落地再续", "降低启动风险"),
        ("媒体支持", "原创3篇约5万+主流宣发等", "可选；打包示意12万，传播≥300万", "按园区媒体需求选配"),
        ("挂牌", "10万/项可选，最多5项", "仍10万/项；本期建议至少1项", "明确启动挂牌意向"),
        ("超额奖励", "出租率≥90% → 20万", "不按整盘90%粗放计提；按贡献细化", "适配非独家现实"),
        ("政策服务", "多作为附赠/陪跑表述", "正式单列并设成果激励", "从附赠变为可结算模块"),
    ]
    for i, r in enumerate(cmp_rows, start=4):
        put(ws, i, list(r))

    wb.save(OUT_CHG_XLS)
    print(f"Wrote {OUT_CHG_XLS}")


def main():
    build_quote_pptx()
    build_change_pptx()
    build_change_xlsx()


if __name__ == "__main__":
    main()
