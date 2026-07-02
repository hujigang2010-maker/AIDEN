"""元谷项目 · 谈判集中汇报 PPT (阶段性集中交付).

把前面所有阶段性成果 (全盘合资公司方案 / 4#5# 楼 2 万方专项 / 上海市科企联
合作 / IP+AI 双轨招商策略 / 排期 / 招商材料) 整合成一份美观、商务、用于今天
谈判沟通的汇报稿。

设计系统:
  - 深海军蓝 (PRIMARY) + 暖橙 (ACCENT) + 金 (GOLD) 商务配色
  - 统一的封面 / 章节分隔页 / KPI 卡片 / 表格 / 时间轴 / 双栏布局
  - 每页统一页眉色条 + 页脚 + 页码
"""
from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

OUT = Path(__file__).with_name("元谷项目谈判集中汇报.pptx")

# ===== 设计系统 配色 =====
PRIMARY = RGBColor(0x0F, 0x24, 0x4E)      # 深海军蓝
PRIMARY2 = RGBColor(0x1B, 0x3A, 0x6B)     # 次蓝
ACCENT = RGBColor(0xF2, 0x7E, 0x2D)       # 暖橙
GOLD = RGBColor(0xC9, 0xA2, 0x4B)         # 金
TEAL = RGBColor(0x18, 0x8B, 0x8B)         # 青
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x1E, 0x8E, 0x5A)
LIGHT = RGBColor(0xF4, 0xF6, 0xFA)        # 浅底
CARD = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD8, 0xDF, 0xEA)
TEXT = RGBColor(0x21, 0x2B, 0x42)
MUTED = RGBColor(0x66, 0x70, 0x86)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

_page = {"n": 0}


def _font(run, size, *, bold=False, color=TEXT, name="微软雅黑"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    from pptx.oxml.ns import qn
    rpr = run._r.get_or_add_rPr()
    rfonts = rpr.find(qn('a:latin'))
    if rfonts is None:
        rfonts = rpr.makeelement(qn('a:latin'), {}); rpr.append(rfonts)
    rfonts.set('typeface', name)
    ea = rpr.find(qn('a:ea'))
    if ea is None:
        ea = rpr.makeelement(qn('a:ea'), {}); rpr.append(ea)
    ea.set('typeface', name)


def _rect(slide, l, t, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE, line_w=1.0):
    sp = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def _text(slide, l, t, w, h, lines, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    """lines: list of (text, size, bold, color) or (text, size, bold, color, space_after)"""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, spec in enumerate(lines):
        text, size, bold, color = spec[0], spec[1], spec[2], spec[3]
        sa = spec[4] if len(spec) > 4 else 3
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(sa)
        p.space_before = Pt(0)
        r = p.add_run(); r.text = text
        _font(r, size, bold=bold, color=color)
    return tb


def bg(slide, color=LIGHT):
    _rect(slide, -0.1, -0.1, 13.6, 7.7, fill=color)


def header(slide, kicker, title, *, accent=ACCENT):
    """统一页眉: 左侧竖色条 + kicker + 标题, 右上角小色块"""
    _rect(slide, 0, 0, 13.333, 1.15, fill=PRIMARY)
    _rect(slide, 0, 1.15, 13.333, 0.06, fill=accent)
    _rect(slide, 0.55, 0.26, 0.12, 0.62, fill=accent)
    _text(slide, 0.85, 0.20, 11.5, 0.34, [(kicker, 11, False, RGBColor(0xB9, 0xC6, 0xDC))])
    _text(slide, 0.85, 0.46, 11.6, 0.6, [(title, 23, True, WHITE)])
    # 右上 logo 区
    _text(slide, 10.6, 0.20, 2.2, 0.8, [("元谷 YUANGU", 12, True, RGBColor(0xE8, 0xC8, 0x7A)),
                                        ("森马 × 胡教授团队", 8.5, False, RGBColor(0xB9, 0xC6, 0xDC))],
          align=PP_ALIGN.RIGHT)


def footer(slide):
    _page["n"] += 1
    _rect(slide, 0, 7.28, 13.333, 0.22, fill=PRIMARY)
    _text(slide, 0.55, 7.28, 9.0, 0.22, [("元谷项目招商运营合作 · 谈判沟通汇报 (阶段性集中交付)", 8, False, RGBColor(0xC6, 0xD0, 0xE2))], anchor=MSO_ANCHOR.MIDDLE)
    _text(slide, 11.8, 7.28, 1.0, 0.22, [(f"{_page['n']:02d}", 8, True, RGBColor(0xE8, 0xC8, 0x7A))], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def kpi_card(slide, l, t, w, h, label, value, sub, *, accent=ACCENT):
    _rect(slide, l, t, w, h, fill=CARD, line=LINE, line_w=0.75)
    _rect(slide, l, t, 0.1, h, fill=accent)
    _text(slide, l + 0.25, t + 0.16, w - 0.4, 0.3, [(label, 10.5, False, MUTED)])
    _text(slide, l + 0.25, t + 0.46, w - 0.4, 0.5, [(value, 21, True, PRIMARY)])
    _text(slide, l + 0.25, t + h - 0.42, w - 0.4, 0.36, [(sub, 9.5, False, MUTED)])


def chip(slide, l, t, w, h, text, *, fill=PRIMARY, fg=WHITE, size=11, bold=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = _rect(slide, l, t, w, h, fill=fill, shape=shape)
    tf = sp.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; _font(r, size, bold=bold, color=fg)
    return sp


def table(slide, l, t, w, headers, rows, *, col_w=None, header_fill=PRIMARY,
          row_h=0.4, head_h=0.42, fs=10.5, hfs=11, zebra=True, accent_first=False):
    n = len(rows) + 1
    cols = len(headers)
    total_h = head_h + row_h * len(rows)
    gt = slide.shapes.add_table(n, cols, Inches(l), Inches(t), Inches(w), Inches(total_h)).table
    # disable default style banding via first/last
    gt.first_row = False; gt.horz_banding = False
    if col_w:
        for j, cw in enumerate(col_w):
            gt.columns[j].width = Inches(cw)
    gt.rows[0].height = Inches(head_h)
    for j, htxt in enumerate(headers):
        c = gt.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = header_fill
        c.margin_left = Inches(0.1); c.margin_right = Inches(0.06)
        c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = htxt; _font(r, hfs, bold=True, color=WHITE)
    for i, row in enumerate(rows, start=1):
        gt.rows[i].height = Inches(row_h)
        for j, val in enumerate(row):
            c = gt.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = (RGBColor(0xED, 0xF1, 0xF8) if (zebra and i % 2 == 0) else CARD)
            c.margin_left = Inches(0.1); c.margin_right = Inches(0.06)
            c.margin_top = Inches(0.01); c.margin_bottom = Inches(0.01)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            is_accent = accent_first and j == 0
            _font(r, fs, bold=is_accent, color=(PRIMARY if is_accent else TEXT))
    return gt


def bullets(slide, l, t, w, h, items, *, fs=14, gap=6, color=TEXT, marker="—", marker_color=ACCENT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_top = 0
    for i, it in enumerate(items):
        lvl = 0; txt = it
        if isinstance(it, tuple):
            txt, lvl = it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.space_before = Pt(0)
        if lvl == 0:
            rm = p.add_run(); rm.text = marker + "  "; _font(rm, fs, bold=True, color=marker_color)
            r = p.add_run(); r.text = txt; _font(r, fs, bold=False, color=color)
        else:
            rm = p.add_run(); rm.text = "      ·  "; _font(rm, fs - 1, bold=False, color=MUTED)
            r = p.add_run(); r.text = txt; _font(r, fs - 1.5, bold=False, color=MUTED)
    return tb


def section_divider(prs, no, zh, en, points):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, -0.1, -0.1, 13.6, 7.7, fill=PRIMARY)
    _rect(s, 0, 0, 4.6, 7.5, fill=PRIMARY2)
    _rect(s, 0.9, 2.0, 0.9, 0.9, fill=ACCENT, shape=MSO_SHAPE.OVAL)
    _text(s, 0.9, 2.0, 0.9, 0.9, [(no, 26, True, WHITE)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _text(s, 0.9, 3.15, 3.4, 1.6, [(zh, 30, True, WHITE), (en, 12, False, RGBColor(0x9F, 0xB2, 0xD0), 0)])
    # 右侧要点
    yy = 2.0
    for pt in points:
        _rect(s, 5.4, yy + 0.12, 0.16, 0.16, fill=ACCENT, shape=MSO_SHAPE.OVAL)
        _text(s, 5.8, yy, 6.8, 0.6, [(pt, 14.5, False, RGBColor(0xE5, 0xEB, 0xF5))])
        yy += 0.72
    return s


def main():
    prs = Presentation()
    prs.slide_width = EMU_W; prs.slide_height = EMU_H
    blank = prs.slide_layouts[6]

    # =================================================================
    # P1 封面
    # =================================================================
    s = prs.slides.add_slide(blank)
    _rect(s, -0.1, -0.1, 13.6, 7.7, fill=PRIMARY)
    _rect(s, 0, 0, 13.333, 0.18, fill=ACCENT)
    _rect(s, 0, 7.32, 13.333, 0.18, fill=GOLD)
    # 装饰斜块
    _rect(s, 9.3, -0.1, 4.3, 7.7, fill=PRIMARY2)
    _rect(s, 11.0, 1.2, 1.7, 1.7, fill=ACCENT, shape=MSO_SHAPE.OVAL)
    _rect(s, 10.2, 4.4, 1.0, 1.0, fill=GOLD, shape=MSO_SHAPE.OVAL)
    _rect(s, 11.9, 5.3, 0.7, 0.7, fill=TEAL, shape=MSO_SHAPE.OVAL)

    _rect(s, 0.9, 2.35, 0.16, 2.1, fill=ACCENT)
    _text(s, 1.25, 2.25, 8.0, 0.5, [("SEMIR 森马 (上海) 国际运营中心 · 元谷项目", 14, False, RGBColor(0xB9, 0xC6, 0xDC))])
    _text(s, 1.2, 2.7, 8.2, 2.2, [
        ("元谷项目招商运营合作", 40, True, WHITE),
        ("谈 判 集 中 汇 报", 30, True, RGBColor(0xE8, 0xC8, 0x7A), 6),
    ])
    _text(s, 1.25, 4.95, 8.0, 0.5, [("IP + AI 双轨  ·  2 万方产业研发办公  ·  阶段性集中交付", 15, False, RGBColor(0xCF, 0xD8, 0xE8))])
    # 底部信息条
    chip(s, 1.25, 5.75, 3.0, 0.5, "胡教授团队 × 森马集团", fill=ACCENT, size=12)
    chip(s, 4.45, 5.75, 2.4, 0.5, "v1.0 谈判稿", fill=PRIMARY2, fg=RGBColor(0xE8,0xC8,0x7A), size=12)
    _text(s, 1.25, 6.5, 8.0, 0.4, [("2026 年 6 月  ·  上海 闵行 大零号湾  ·  本汇报仅用于商务谈判沟通", 10.5, False, RGBColor(0x9F, 0xB2, 0xD0))])

    # =================================================================
    # P2 汇报议程
    # =================================================================
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "AGENDA", "汇报议程 · 七大板块")
    agenda = [
        ("01", "项目与共识", "战略定位 / 双重身份 / 本次谈判要点"),
        ("02", "我们能带来什么", "六大稀缺资源资产化"),
        ("03", "2 万方招商方案", "IP+AI 双轨 / 四级漏斗 / 六大出彩点"),
        ("04", "节点与排期", "9/30 与 5/1 两个硬节点 / 月度签约率"),
        ("05", "品牌与活动", "5 项挂牌 / 6 场沙龙 / 5·22 峰会借势"),
        ("06", "商业条款与测算", "月费 / 佣金 / 挂牌 / 沙龙 / ROI"),
        ("07", "投决建议与下一步", "三问决策 / 30 天启动"),
    ]
    x0, y0, cw, ch, gx, gy = 0.7, 1.5, 6.05, 1.18, 0.25, 0.2
    for i, (no, zh, en) in enumerate(agenda):
        col = i % 2; rowi = i // 2
        l = x0 + col * (cw + gx); t = y0 + rowi * (ch + gy)
        _rect(s, l, t, cw, ch, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, l, t, 1.1, ch, fill=PRIMARY)
        _text(s, l, t, 1.1, ch, [(no, 24, True, RGBColor(0xE8, 0xC8, 0x7A))], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, l + 1.3, t + 0.22, cw - 1.45, 0.5, [(zh, 15.5, True, PRIMARY)])
        _text(s, l + 1.3, t + 0.68, cw - 1.45, 0.45, [(en, 10.5, False, MUTED)])
    # 第 7 项 (01..07 共 7 个, 居中放最后一行)
    footer(s)

    # =================================================================
    # 章节 01
    # =================================================================
    section_divider(prs, "01", "项目与共识", "PROJECT & CONSENSUS", [
        "元谷:大零号湾文创融合核心区 + 上海唯一科技时尚特色小镇",
        "本期聚焦 4#+5# 楼 5F+ 共约 2 万方产业研发办公",
        "战略已对齐:从纯二次元升级为 IP + AI 双轨",
        "本次谈判:锁定合作模式 + 商业条款 + 启动时间",
    ])

    # P 项目定位
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "01 · 项目与共识", "元谷的双重战略身份与产业底盘")
    kpi_card(s, 0.7, 1.5, 2.95, 1.5, "总建筑面积", "22 万㎡", "1#-6# 共 6 栋")
    kpi_card(s, 3.83, 1.5, 2.95, 1.5, "本期招商范围", "2 万㎡", "4#+5# 楼 5F+ 产业办公", accent=GOLD)
    kpi_card(s, 6.96, 1.5, 2.95, 1.5, "TOD 单日客流", "5-7 万人次", "15 号线元江路站", accent=TEAL)
    kpi_card(s, 10.09, 1.5, 2.55, 1.5, "15min 覆盖", "24万+12万", "居住 + 产业人口", accent=PRIMARY2)
    _text(s, 0.7, 3.3, 12, 0.4, [("两大稀缺政策身份", 15, True, PRIMARY)])
    chip(s, 0.7, 3.78, 5.9, 0.62, "大零号湾文创融合核心区 (闵行五大中心之一)", fill=PRIMARY, size=12.5)
    chip(s, 6.75, 3.78, 5.9, 0.62, "上海市唯一科技时尚特色小镇", fill=ACCENT, size=12.5)
    bullets(s, 0.7, 4.7, 12.0, 2.2, [
        "能级比肩漕河泾、张江;享区级最高 1,000 万 + 市级最高 5,000 万产业专项",
        "产业愿景:构建『时尚研发』+『文创转化』双向赋能标杆, 打造 AI+IP 双赛道协同园区",
        "森马给底盘 + 危总给地利 + 胡教授给打法 → 三方协同, 优势互补",
    ], fs=14)
    footer(s)

    # P 楼宇底盘
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "01 · 项目与共识", "楼宇产业底盘:本期聚焦 4# / 5# 楼 5F+")
    table(s, 0.7, 1.55, 11.95,
          ["楼栋", "1-4F 业态", "5F 及以上", "本期定位"],
          [
              ["1# 楼", "零售", "森马集团总部办公", "—"],
              ["2# 楼", "二次元主题 Livehouse / 秀场", "—", "活动主场"],
              ["3# 楼", "休闲运动 / 萌宠空间", "酒店", "配套"],
              ["★ 4# 楼", "潮玩艺术中心 / 4F 直播中心", "潮玩产业集群 ≈ 1 万㎡", "AI 主轴 · 国际创意层"],
              ["★ 5# 楼", "动漫书店 / 休闲娱乐", "潮玩产业集群 ≈ 1 万㎡", "IP 主轴 · 产业总部层"],
              ["6# 楼", "品质生活 / 特色餐饮 / 商务宴请", "—", "配套"],
          ],
          col_w=[1.5, 4.7, 3.45, 2.3], row_h=0.62, accent_first=True)
    _text(s, 0.7, 6.35, 12, 0.5, [("★ 4# 楼 + 5# 楼 5F+ 共约 2 万㎡是元谷的『产业心脏』;商业 5.2 万方由森马商业部另行招商, 非本期重点。", 12, True, ACCENT)])
    footer(s)

    # =================================================================
    # 章节 02
    # =================================================================
    section_divider(prs, "02", "我们能带来什么", "WHAT WE BRING", [
        "六大稀缺资源, 全部可资产化为招商工具",
        "北欧会客厅 / 福布斯 / 科技开放麦 / AI 腾讯",
        "仲量联行爬楼大数据 (已投入 2.6 万) / 追觅基金",
        "复旦住房政策研究中心 + 上海市科技企业联合会",
    ])

    s = prs.slides.add_slide(blank); bg(s)
    header(s, "02 · 我们能带来什么", "六大稀缺资源 → 招商引擎")
    res = [
        ("北欧创新国际会客厅", "海外 IP / 技术进入中国首站", "4# 楼挂牌", ACCENT),
        ("福布斯产业影响力奖", "国际品牌势能 + 年度评选 IP", "元谷专场", GOLD),
        ("科技开放麦 + 6 场沙龙", "活动即招商, 单场 ≥30 客户", "年度 IP 化", TEAL),
        ("AI 腾讯 + 算力补贴", "新注册 AI 公司 3 月免费 / 算力 85 折", "AI 招商抓手", PRIMARY2),
        ("仲量联行爬楼大数据", "200 家精准客户清单, 转化率 +30%", "已投入 2.6 万", ACCENT),
        ("追觅基金 + CVC 创投", "返投落地, 资本驱动签约", "元谷产业基金", GOLD),
    ]
    x0, y0, cw, ch, gx, gy = 0.7, 1.5, 3.85, 1.62, 0.18, 0.2
    for i, (t1, t2, tag, ac) in enumerate(res):
        col = i % 3; rowi = i // 3
        l = x0 + col * (cw + gx); t = y0 + rowi * (ch + gy)
        _rect(s, l, t, cw, ch, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, l, t, cw, 0.1, fill=ac)
        _text(s, l + 0.22, t + 0.24, cw - 0.4, 0.5, [(t1, 13.5, True, PRIMARY)])
        _text(s, l + 0.22, t + 0.74, cw - 0.4, 0.6, [(t2, 10.5, False, MUTED)])
        chip(s, l + 0.22, t + ch - 0.46, 1.9, 0.32, tag, fill=ac, size=9.5)
    _text(s, 0.7, 6.9, 12, 0.4, [("两大政府/学术平台已确认:复旦大学住房政策研究中心 · 元谷分中心  +  上海市科技企业联合会 · 元谷产业基地", 12, True, PRIMARY)])
    footer(s)

    # =================================================================
    # 章节 03
    # =================================================================
    section_divider(prs, "03", "2 万方招商方案", "LEASING STRATEGY", [
        "战略转向:纯二次元 → IP + AI 双轨",
        "客户策略:优先 2,000-20,000㎡ 大客户 (转化率 110:1)",
        "四级招商漏斗 + 招商六大『出彩点』",
        "五档客户配比, 与森马原产业规划完全对齐",
    ])

    # P 战略转向
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "03 · 2 万方招商方案", "核心转向:从『纯二次元』→『IP + AI 双轨』")
    table(s, 0.7, 1.5, 11.95,
          ["维度", "原方向 (纯二次元)", "新方向 (IP + AI 双轨)"],
          [
              ["招商面", "200-500㎡ 小户, 招满需 400-800 家", "AI/IP 中大型, 招满仅需 30-60 家"],
              ["支付能力", "单家年租 < 30 万", "单家年租 100-500 万 (真金白银)"],
              ["政府政策", "拿政策难, 创新基金不投", "享高新 + 专精特新 + 创新券 + 闵行专项"],
              ["与森马协同", "二次元与商业难协同", "IP+AI 与森马时尚科技定位契合"],
              ["先例", "无 2 万方成功案例", "杨浦 1 万方『AI+IP』已通过先例"],
          ],
          col_w=[1.8, 4.85, 5.3], row_h=0.6, accent_first=True)
    chip(s, 0.7, 5.0, 11.95, 0.7, "结论:主轴改 AI 拿政府资源 + 资本, 保留 IP 概念 (视频/漫画/直播/动漫全链条) 兼顾威总文化导向", fill=ACCENT, size=13)
    bullets(s, 0.7, 5.9, 12, 1.1, [
        "4# 楼 5F+ = AI 主轴 (国际创意层);5# 楼 5F+ = IP 主轴 (产业总部层)",
        "行业转化率新基准 110:1 → 招 1 家大户 ≫ 招 10 家小户, 效率差 5 倍",
    ], fs=13)
    footer(s)

    # P 四级漏斗
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "03 · 2 万方招商方案", "四级招商漏斗 + 五档客户配比")
    funnel = [
        ("L1 牌照锚定", "头部央企 / 行业协会", "AI 潮玩产业基地 + 潮玩次元专委 双牌照", "6,000㎡", ACCENT),
        ("L2 资本招商", "中型 AI / IP 企业", "追觅基金 + AI 腾讯 + CVC 创投", "5,000㎡", GOLD),
        ("L3 大数据爬楼", "小型潮玩 / 服务机构", "仲量联行爬楼大数据 (已购)", "6,000-8,000㎡", TEAL),
        ("L4 活动带流", "中小型服务机构", "6 场沙龙 + 5·22 峰会 + 福布斯", "3,000-4,000㎡", PRIMARY2),
    ]
    yy = 1.55
    for name, who, how, area, ac in funnel:
        _rect(s, 0.7, yy, 11.95, 1.12, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, 0.7, yy, 2.5, 1.12, fill=ac)
        _text(s, 0.7, yy, 2.5, 1.12, [(name, 15, True, WHITE)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, 3.4, yy + 0.16, 5.6, 0.4, [(who, 13, True, PRIMARY)])
        _text(s, 3.4, yy + 0.58, 6.6, 0.45, [(how, 11, False, MUTED)])
        chip(s, 10.4, yy + 0.33, 2.0, 0.46, area, fill=ac, size=12)
        yy += 1.24
    _text(s, 0.7, 6.6, 12, 0.5, [("★ 漏斗合计可产生 20,000-23,000㎡ 招商管道 → 实际签约 18,000-20,000㎡ (满租);五档配比 10%央企+10%配套+20%中型+20%服务+40%小型, 与森马原规划一致", 11, True, ACCENT)])
    footer(s)

    # P 六大出彩点
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "03 · 2 万方招商方案", "招商六大『出彩点』(VS 普通园区)")
    pts = [
        ("① 牌照即招商", "双牌照前置, 客户『送进来』不是『拉进来』"),
        ("② 基金即招商", "追觅 1:1.5 返投 + CVC + 元谷产业基金"),
        ("③ 数据即招商", "仲量联行爬楼数据, 转化率 +30%"),
        ("④ 峰会即招商", "5·22 峰会 200+ VIP, 1 天锁 30-50 客户"),
        ("⑤ 沙龙即招商", "6 场 × ≥30 客户 = 5-12 家直接成果"),
        ("⑥ 学术即招商", "复旦 + 上海交大 + 北大 三校背书"),
    ]
    x0, y0, cw, ch, gx, gy = 0.7, 1.55, 5.9, 1.42, 0.15, 0.2
    for i, (t1, t2) in enumerate(pts):
        col = i % 2; rowi = i // 2
        l = x0 + col * (cw + gx); t = y0 + rowi * (ch + gy)
        _rect(s, l, t, cw, ch, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, l, t, 0.1, ch, fill=ACCENT)
        _text(s, l + 0.28, t + 0.2, cw - 0.5, 0.5, [(t1, 15, True, PRIMARY)])
        _text(s, l + 0.28, t + 0.74, cw - 0.5, 0.55, [(t2, 11.5, False, MUTED)])
    _text(s, 0.7, 6.95, 12, 0.4, [("普通园区只做 ③④;元谷六维齐发 → 招商速度领先 6-12 个月, 单方租金溢价 0.2-0.4 元/㎡/天", 12, True, ACCENT)])
    footer(s)

    # =================================================================
    # 章节 04
    # =================================================================
    section_divider(prs, "04", "节点与排期", "MILESTONES & SCHEDULE", [
        "硬节点一:2026/9/30 完成 2,000㎡ 签约 (T+100 天)",
        "硬节点二:2027/5/1 项目开业, 50%+ 签约 (T+314 天)",
        "全周期:2028/5/1 达成 90%+ 满租",
        "五阶段倒推 + 月度签约率追踪",
    ])

    # P 两个硬节点 + 五阶段
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "04 · 节点与排期", "两个硬节点 · 倒推全部排期")
    kpi_card(s, 0.7, 1.5, 5.9, 1.45, "硬节点一 · 2026/9/30 (T+100 天)", "2,000㎡ 签约", "可含直播基地 / 共享设计中心擦边球", accent=RED)
    kpi_card(s, 6.75, 1.5, 5.9, 1.45, "硬节点二 · 2027/5/1 (T+314 天)", "开业 + 50% 签约", "2 万方达成 50% + 消费氛围成型", accent=ACCENT)
    table(s, 0.7, 3.2, 11.95,
          ["阶段", "时间窗口", "累计签约目标", "里程碑"],
          [
              ["阶段 0 准备", "6/22 - 7/15", "0", "团队就位 + 资源接入 + 物料"],
              ["阶段 1 抢节点", "7/15 - 9/30", "2,300㎡ (含擦边球)", "★ 9/30 硬节点达成"],
              ["阶段 2 加速", "10/1 - 12/31", "5,500㎡ (27%)", "L3 爬楼 + 沙龙 + 挂牌 ③④"],
              ["阶段 3 开业", "1/1 - 5/1", "10,400㎡ (52%)", "★ 5/1 开业典礼"],
              ["阶段 4 满租", "5/2 - 次年 5/1", "18,000㎡+ (90%)", "沙龙 IP 化 + 续约"],
          ],
          col_w=[2.0, 2.6, 3.35, 4.0], row_h=0.52, accent_first=True)
    footer(s)

    # P 月度签约率
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "04 · 节点与排期", "月度签约率推进 (按 803 元/㎡ 年租金)")
    table(s, 0.7, 1.5, 11.95,
          ["月份", "本月新增", "累计签约", "签约率", "累计年租金", "里程碑"],
          [
              ["2026/7", "300㎡", "300㎡", "1.5%", "24 万", "战队就位 + 物料"],
              ["2026/8", "800㎡", "1,100㎡", "5.5%", "88 万", "L1 牌照 + 大客户"],
              ["2026/9", "1,200㎡", "2,300㎡", "11.5%", "185 万", "★ 9/30 节点达成"],
              ["2026/12", "—", "5,500㎡", "27.5%", "442 万", "Q4 攻坚收官"],
              ["2027/4", "1,500㎡", "10,400㎡", "52%", "836 万", "★ 5/1 节点达成"],
              ["2027/5", "1,500㎡", "11,900㎡", "59.5%", "956 万", "开业典礼"],
          ],
          col_w=[1.9, 1.85, 2.0, 1.6, 2.3, 2.3], row_h=0.5, accent_first=True)
    chip(s, 0.7, 5.3, 11.95, 0.66, "9/30 节点 2,000㎡ 来源:大客户 1,000-2,000㎡ + 直播基地 800-1,000㎡ + AI 设计中心 500-800㎡ + 服务中心 300-500㎡ → 中性期望 3,200-5,800㎡, 达成概率 >95%", fill=PRIMARY, size=11.5)
    _text(s, 0.7, 6.2, 12, 0.5, [("单 12 个月口径:甲方租金累计回收 ≈ 956 万 + 物业 ≈ 90 万 = 1,046 万元 → 已现金回正", 12.5, True, GREEN)])
    footer(s)

    # =================================================================
    # 章节 05
    # =================================================================
    section_divider(prs, "05", "品牌与活动", "BRAND & EVENTS", [
        "5 项产业牌照, 一次性挂牌奖励 30 万/项",
        "6 场产业沙龙, 每场 ≥30 个目标产业客户",
        "5·22 AI 商业化峰会借势, 200+ VIP 一日转化",
        "活动 IP 化 → 年度品牌势能 + 媒体声量 5 亿次",
    ])

    # P 五项挂牌
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "05 · 品牌与活动", "5 项产业牌照 (5 × 30 万 = 150 万一次性激励)")
    bands = [
        ("①", "AI 潮玩产业基地", "中国动漫集团", "T+3 月"),
        ("②", "潮玩次元商业专委会", "中国百货商业协会", "T+3 月"),
        ("③", "复旦大学住房政策研究中心 · 元谷分中心", "复旦大学", "T+9 月"),
        ("④", "上海市科技企业联合会 · 元谷产业基地", "上海市科企联", "T+6 月"),
        ("⑤", "福布斯产业影响力奖 · 元谷专场", "福布斯", "T+12 月"),
    ]
    yy = 1.55
    for no, name, body, node in bands:
        _rect(s, 0.7, yy, 11.95, 0.86, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, 0.7, yy, 0.86, 0.86, fill=PRIMARY)
        _text(s, 0.7, yy, 0.86, 0.86, [(no, 22, True, RGBColor(0xE8, 0xC8, 0x7A))], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, 1.75, yy + 0.16, 7.4, 0.55, [(name, 14, True, PRIMARY)])
        _text(s, 1.75, yy + 0.5, 7.4, 0.32, [("出牌方:" + body, 10.5, False, MUTED)])
        chip(s, 9.4, yy + 0.24, 1.4, 0.38, node, fill=TEAL, size=11)
        chip(s, 10.95, yy + 0.24, 1.55, 0.38, "30 万", fill=ACCENT, size=12)
        yy += 0.98
    footer(s)

    # P 6 场沙龙 + 5·22
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "05 · 品牌与活动", "6 场产业沙龙 (每场 ≥30 客户) + 5·22 峰会借势")
    table(s, 0.7, 1.5, 7.55,
          ["#", "时间", "主题", "联办方"],
          [
              ["1", "T+1月", "AI+潮玩 (借势5·22)", "中动漫+腾讯"],
              ["2", "T+3月", "潮玩出海", "北欧+福布斯"],
              ["3", "T+5月", "投融资路演", "追觅+招行+长江"],
              ["4", "T+7月", "设计与创意", "上海交大+科企联"],
              ["5", "T+9月", "内容IP·Z世代", "中百协+中动漫"],
              ["6", "T+11月", "政策补贴·小镇", "闵行科协+复旦"],
          ],
          col_w=[0.5, 1.2, 2.85, 3.0], row_h=0.5, fs=10, accent_first=True)
    # 右侧 5·22 借势卡
    _rect(s, 8.55, 1.5, 4.1, 4.5, fill=PRIMARY, line=None)
    _rect(s, 8.55, 1.5, 4.1, 0.6, fill=ACCENT)
    _text(s, 8.75, 1.5, 3.8, 0.6, [("5·22 AI 商业化峰会借势", 13, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, 8.78, 2.3, 3.7, 3.6, [
        "200+ VIP (北大/复旦/中行/长江/招行/金浦/铂帝)",
        "元谷设专属展位 + 闭门 1V1 招商台",
        "增设『AI 潮玩产业影响力榜』+ 福布斯背书",
        "峰会嘉宾 → 转化为沙龙 #1 主题嘉宾",
        "闭门晚宴 → 元谷创始合伙人圈",
    ], fs=11.5, gap=10, color=RGBColor(0xE5, 0xEB, 0xF5), marker="▸", marker_color=RGBColor(0xE8, 0xC8, 0x7A))
    _text(s, 0.7, 6.3, 12, 0.5, [("每场沙龙触达 ≥30 家产业客户 + ≥100 万次媒体曝光;年度活动 IP 化 → 媒体声量 5 亿次/年", 12, True, ACCENT)])
    footer(s)

    # =================================================================
    # 章节 06
    # =================================================================
    section_divider(prs, "06", "商业条款与测算", "TERMS & ROI", [
        "基础月费:推荐 12 万/月 (1-3 人配置可选)",
        "招商佣金:1.5-2 个月年租金 (按面积阶梯)",
        "挂牌 30 万/项 × 5 + 沙龙 5 万/场 × 6",
        "24 月双向账本:ROI 1:3.88",
    ])

    # P 商业条款
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "06 · 商业条款与测算", "向森马提议的合作条款 (五项收入)")
    table(s, 0.7, 1.5, 11.95,
          ["收入类别", "推荐方案", "支付方", "结算"],
          [
              ["基础月费 Retainer", "12 万元/月 (区间 6-18 万, 1/2/3 人)", "森马", "按月预付"],
              ["招商佣金", "成交年租金 1.5 / 1.75 / 2.0 个月 (按面积)", "森马", "起租后 30 日"],
              ["挂牌奖励", "30 万元/项 × 5 项 = 150 万 (一次性)", "森马", "挂牌后 30 日"],
              ["沙龙执行费", "5 万元/场 × 6 场 = 30 万 + 净利 30/70 分润", "森马 + 外部", "按场"],
              ["超额奖励", "24 月满租率 ≥95% → 100 万一次性", "森马", "达成后"],
          ],
          col_w=[2.6, 5.55, 1.7, 2.1], row_h=0.58, accent_first=True)
    bullets(s, 0.7, 5.05, 12, 1.6, [
        "团队配置:2 人驻场 + CSO 顾问 (胡教授每周 ≥2 天), 月度刚性成本约 9.5 万 → 推荐月费 12 万",
        "客户归属保护期 24 个月, 甲方不得绕开签约;返投基金客户佣金额外 +0.25 个月",
        "服务期 24 个月 (首期), 满租后乙方享优先续约权",
    ], fs=13)
    footer(s)

    # P ROI 双向账本
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "06 · 商业条款与测算", "24 月双向账本:森马『花一份钱赚四份钱』")
    kpi_card(s, 0.7, 1.55, 3.85, 1.7, "甲方付出 (24 月)", "714 万元", "月费+佣金+挂牌+沙龙", accent=RED)
    kpi_card(s, 4.75, 1.55, 3.85, 1.7, "甲方收入 (24 月)", "2,769 万元", "租金 + 物业", accent=GREEN)
    kpi_card(s, 8.8, 1.55, 3.85, 1.7, "投入产出比 ROI", "1 : 3.88", "尚不含资产增值", accent=ACCENT)
    table(s, 0.7, 3.6, 11.95,
          ["维度", "金额 / 指标", "说明"],
          [
              ["12 个月投入", "≈ 684 万元", "月费+佣金+挂牌+沙龙+开业+硬件"],
              ["12 个月收入", "≈ 1,046 万元", "租金 956 万 + 物业 90 万"],
              ["12 个月净收益", "≈ +362 万元", "★ 单 12 月已现金回正"],
              ["永续年化租金", "1,606 万元/年", "2 万㎡ 满租 × 2.2 元/㎡/天"],
              ["资产估值增量", "≈ 2.47 亿元", "按 8% 折现 (租金+物业+配套)"],
          ],
          col_w=[2.8, 3.0, 6.15], row_h=0.5, accent_first=True)
    footer(s)

    # =================================================================
    # 章节 07
    # =================================================================
    section_divider(prs, "07", "投决建议与下一步", "DECISION & NEXT STEP", [
        "三问决策, 建议本周内确定",
        "30 天内启动, 赶 9/30 硬节点",
        "错过 7 月窗口, 节点风险陡增",
        "先签 2 万方专项 → 12 月成势后承接全盘",
    ])

    # P 投决建议
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "07 · 投决建议与下一步", "三问决策 (建议本周内确定)")
    qs = [
        ("Q1", "是否接受 IP+AI 双轨战略调整?", "建议:是 (杨浦区已有 1 万方先例)"),
        ("Q2", "是否接受 12 万月费 + 1.5-2 月佣金 + 5 项挂牌 + 6 场沙龙条款?", "建议:是 (ROI 1:3.88, 12 月即回正)"),
        ("Q3", "是否同意 7 月初启动 (赶 9/30 硬节点)?", "建议:是 (错过 7 月窗口风险陡增)"),
    ]
    yy = 1.6
    for no, q, a in qs:
        _rect(s, 0.7, yy, 11.95, 1.18, fill=CARD, line=LINE, line_w=0.75)
        _rect(s, 0.7, yy, 1.1, 1.18, fill=PRIMARY)
        _text(s, 0.7, yy, 1.1, 1.18, [(no, 20, True, RGBColor(0xE8, 0xC8, 0x7A))], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, 2.0, yy + 0.2, 10.4, 0.55, [(q, 14, True, PRIMARY)])
        _text(s, 2.0, yy + 0.72, 10.4, 0.4, [(a, 12, False, ACCENT)])
        yy += 1.32
    _text(s, 0.7, 6.55, 12, 0.5, [("分步策略:先签 4#+5# 楼 2 万方专项试点 → 12 个月成势后, 用全盘方案承接 5.2 万方 + 合资公司", 12.5, True, PRIMARY)])
    footer(s)

    # P 30 天启动路线
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "07 · 投决建议与下一步", "30 / 60 / 90 天启动路线图")
    steps = [
        ("30 天", "MOU 签订 + 2-3 人战队就位 + 仲量联行数据接入 + 5·22 峰会借势", ACCENT),
        ("60 天", "沙龙 #1 (AI+潮玩) 完成 + AI 潮玩产业基地挂牌筹备 + 签约首单", GOLD),
        ("90 天", "双牌照正式挂牌 + 4,000㎡ 签约 (L1 牌照客户 3 家入驻)", TEAL),
        ("180 天", "沙龙过半 + 科企联挂牌 + 8,000㎡ 签约", PRIMARY2),
    ]
    yy = 1.6
    for i, (d, txt, ac) in enumerate(steps):
        _rect(s, 1.4, yy, 0.16, 1.0 if i < 3 else 0.0, fill=LINE) if i < 3 else None
        _rect(s, 1.05, yy, 0.85, 0.85, fill=ac, shape=MSO_SHAPE.OVAL)
        _text(s, 1.05, yy, 0.85, 0.85, [(d.split()[0], 15, True, WHITE)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _rect(s, 2.2, yy, 10.4, 0.85, fill=CARD, line=LINE, line_w=0.75)
        _text(s, 2.45, yy, 2.0, 0.85, [(d, 14, True, ac)], anchor=MSO_ANCHOR.MIDDLE)
        _text(s, 4.3, yy, 8.1, 0.85, [(txt, 12, False, TEXT)], anchor=MSO_ANCHOR.MIDDLE)
        yy += 1.15
    footer(s)

    # P 附:阶段性交付清单
    s = prs.slides.add_slide(blank); bg(s)
    header(s, "附录 · 阶段性交付", "已完成的全部交付物 (本汇报为集中版)")
    table(s, 0.7, 1.5, 11.95,
          ["模块", "交付物", "形态"],
          [
              ["全盘合资公司方案", "战略汇报 + 收益测算 + 联合运营协议", "PPT / Excel / Word"],
              ["4#+5# 楼 2 万方专项", "九段式方案 + 测算表 + 专项协议", "PPT / Excel / Word"],
              ["上海市科企联合作", "合作服务方案 (服务+报价+中介费逻辑)", "Word / Markdown"],
              ["招商策略与排期", "策略排期 + 排期总表 + 文字执行版", "PPT / Excel / Word"],
              ["招商材料工具包", "招商手册 + 百问话术 + 政策汇编 + 租金对标", "Word ×3 / Excel"],
              ["★ 谈判集中汇报 (本件)", "整合全部交付的谈判用一页流", "PPT (本件)"],
          ],
          col_w=[3.0, 6.35, 2.6], row_h=0.56, accent_first=True)
    _text(s, 0.7, 6.4, 12, 0.5, [("全部源文件与可下载 ZIP 已归档;本 PPT 为今日谈判用『集中浓缩版』, 详情可随时展开任一模块原件。", 11.5, False, MUTED)])
    footer(s)

    # =================================================================
    # 尾页
    # =================================================================
    s = prs.slides.add_slide(blank)
    _rect(s, -0.1, -0.1, 13.6, 7.7, fill=PRIMARY)
    _rect(s, 0, 0, 13.333, 0.18, fill=ACCENT)
    _rect(s, 0, 7.32, 13.333, 0.18, fill=GOLD)
    _rect(s, 1.2, 2.4, 0.16, 2.2, fill=ACCENT)
    _text(s, 1.55, 2.5, 11.0, 2.6, [
        ("100 天 · 2,000 方", 46, True, WHITE),
        ("IP + AI 双轨 · 借势政府 · 闭环招商", 22, False, RGBColor(0xE8, 0xC8, 0x7A), 10),
        ("让元谷成为中国潮玩与 AI 产业走向世界的国际客厅", 16, False, RGBColor(0xCF, 0xD8, 0xE8), 0),
    ])
    chip(s, 1.55, 5.6, 4.0, 0.55, "胡教授团队 × 森马集团", fill=ACCENT, size=13)
    _text(s, 1.55, 6.45, 11, 0.4, [("谈判集中汇报 v1.0  ·  2026 年 6 月  ·  期待与森马达成合作", 12, False, RGBColor(0x9F, 0xB2, 0xD0))])

    prs.save(OUT)
    print(f"Wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
