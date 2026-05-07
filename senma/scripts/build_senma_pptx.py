"""
SENMA · MODE Hub 招商方案 PPT 生成脚本（商务汇报版 · 16:9 · 35 页）

视觉风格：时尚 / 紫红 + 金，区别于冠松的智驾蓝
- 主色板：
    深紫红 #5C1A4E 主基调（时尚 / 高端）
    玫红   #C2185B 强调 / 数据
    金     #C9A24A 重点 / 数字高亮
    云灰   #F4F0F4 正文背景
    炭黑   #1B1F2A 正文字
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pathlib import Path

NAVY = RGBColor(0x5C, 0x1A, 0x4E)     # 深紫红（替代深海蓝）
BLUE = RGBColor(0xC2, 0x18, 0x5B)     # 玫红（替代智驾蓝）
GOLD = RGBColor(0xC9, 0xA2, 0x4A)
CLOUD = RGBColor(0xF4, 0xF0, 0xF4)    # 微暖云灰
INK = RGBColor(0x1B, 0x1F, 0x2A)
GREY = RGBColor(0x6B, 0x73, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xE0, 0xD0, 0xDC)
GREEN = RGBColor(0x2F, 0xA3, 0x6F)
RED = RGBColor(0xD0, 0x4A, 0x4A)
PURPLE = RGBColor(0x6B, 0x4A, 0xC0)

CN_FONT = "WenQuanYi Micro Hei"
EN_FONT = "Calibri"

SLIDES = []


# ----------- 工具函数（与冠松版一致） -----------
def set_run(run, text, *, size=14, bold=False, color=INK, italic=False,
            font_cn=CN_FONT, font_en=EN_FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_en
    rPr = run._r.get_or_add_rPr()
    for tag in ("ea", "cs"):
        existing = rPr.find(qn(f"a:{tag}"))
        if existing is not None:
            rPr.remove(existing)
        ea = rPr.makeelement(qn(f"a:{tag}"), {"typeface": font_cn})
        rPr.append(ea)


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, fill=None, line=None,
             italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if fill is not None:
        box.fill.solid(); box.fill.fore_color.rgb = fill
    else:
        box.fill.background()
    if line is not None:
        box.line.color.rgb = line; box.line.width = Pt(0.5)
    else:
        box.line.fill.background()
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        set_run(run, ln, size=size, bold=bold, color=color, italic=italic)
    return box


def add_rect(slide, x, y, w, h, *, fill=NAVY, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(0.5)
    shp.shadow.inherit = False
    shp.text_frame.text = ""
    return shp


def add_round(slide, x, y, w, h, text="", *, fill=BLUE, color=WHITE,
              size=12, bold=True, align=PP_ALIGN.CENTER):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = align
    set_run(p.add_run(), text, size=size, bold=bold, color=color)
    return shp


def add_table(slide, x, y, w, h, header, rows, *, header_fill=NAVY,
              header_color=WHITE, zebra=(WHITE, CLOUD), header_size=11,
              body_size=10, col_widths=None, body_align=PP_ALIGN.LEFT):
    cols = len(header); n_rows = len(rows) + 1
    ts = slide.shapes.add_table(n_rows, cols, x, y, w, h)
    table = ts.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw
    for j, htxt in enumerate(header):
        cell = table.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Pt(4); tf.margin_right = Pt(4)
        tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
        tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        set_run(p.add_run(), htxt, size=header_size, bold=True, color=header_color)
    for i, row in enumerate(rows, start=1):
        bg = zebra[(i - 1) % 2]
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Pt(4); tf.margin_right = Pt(4)
            tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
            tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = body_align
            set_run(p.add_run(), str(val), size=body_size, color=INK)
    return ts


def add_chrome(slide, prs, *, page_no, phase_label="", page_title="", subtitle=""):
    sw, sh = prs.slide_width, prs.slide_height
    add_rect(slide, 0, 0, sw, Emu(380000), fill=NAVY)
    if phase_label:
        add_round(slide, Inches(0.5), Inches(0.18), Inches(2.0), Inches(0.36),
                  phase_label, fill=GOLD, color=NAVY, size=11, bold=True)
    if page_title:
        add_text(slide, Inches(2.7), Inches(0.10), Inches(8.5), Inches(0.55),
                 page_title, size=22, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(2.7), Inches(0.55), Inches(8.5), Inches(0.30),
                 subtitle, size=11, color=CLOUD, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 0, sh - Emu(80000), sw, Emu(80000), fill=GOLD)
    add_text(slide, Inches(0.5), Inches(7.05), Inches(8.0), Inches(0.30),
             "SENMA · MODE Hub · 森马时尚生态总部园区 · 招商方案 v1.0",
             size=9, color=GREY)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(2.0), Inches(0.30),
             f"{page_no} / 0", size=9, color=GREY, align=PP_ALIGN.RIGHT)


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    SLIDES.append(s)
    return s


def section_cover(prs, idx_label, title_cn, title_en, points):
    s = new_slide(prs)
    sw, sh = prs.slide_width, prs.slide_height
    add_rect(s, 0, 0, sw, sh, fill=NAVY)
    add_rect(s, 0, Inches(3.7), sw, Emu(40000), fill=GOLD)
    add_text(s, Inches(0.8), Inches(1.0), Inches(8), Inches(0.6),
             idx_label, size=18, bold=True, color=GOLD)
    add_text(s, Inches(0.8), Inches(1.7), Inches(11), Inches(1.3),
             title_cn, size=44, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.0), Inches(11), Inches(0.6),
             title_en, size=18, color=CLOUD, italic=True)
    for i, p in enumerate(points):
        y = Inches(4.2) + Inches(0.55) * i
        add_round(s, Inches(0.8), y, Inches(0.4), Inches(0.4), str(i + 1),
                  fill=GOLD, color=NAVY, size=14, bold=True)
        add_text(s, Inches(1.4), y, Inches(11), Inches(0.4), p,
                 size=15, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    return s


def anchor_one_pager(prs, *, idx, label, brand, title, why_anchor, what_lacks,
                     unique_offer, deal_terms, next_steps):
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label=f"Phase 2 · 任务 3 · 提案 {idx}",
               page_title=f"{brand} · 一页纸定制提案",
               subtitle=title)
    top_y = Inches(1.20); top_h = Inches(1.55)
    add_rect(s, Inches(0.5), top_y, Inches(6.15), top_h, fill=NAVY)
    add_text(s, Inches(0.7), top_y + Inches(0.10), Inches(5.8), Inches(0.4),
             f"为什么是 {label}", size=14, bold=True, color=GOLD)
    add_text(s, Inches(0.7), top_y + Inches(0.50), Inches(5.8), Inches(1.0),
             why_anchor, size=11, color=WHITE)
    add_rect(s, Inches(6.85), top_y, Inches(6.15), top_h, fill=BLUE)
    add_text(s, Inches(7.05), top_y + Inches(0.10), Inches(5.8), Inches(0.4),
             "他正缺什么", size=14, bold=True, color=GOLD)
    add_text(s, Inches(7.05), top_y + Inches(0.50), Inches(5.8), Inches(1.0),
             what_lacks, size=11, color=WHITE)
    mid_y = Inches(2.95)
    add_text(s, Inches(0.5), mid_y, Inches(12.5), Inches(0.4),
             "MODE Hub + 森马 · 独家供给", size=14, bold=True, color=NAVY)
    add_rect(s, Inches(0.5), mid_y + Inches(0.40), Inches(12.5), Emu(20000), fill=GOLD)
    cw = Inches(3.10); gx = Inches(0.05); cy = mid_y + Inches(0.55); ch = Inches(1.45)
    colors = [BLUE, NAVY, GOLD, GREEN]
    for i, (t, d) in enumerate(unique_offer):
        x = Inches(0.5) + (cw + gx) * i
        add_rect(s, x, cy, cw, Inches(0.45), fill=colors[i % 4])
        add_text(s, x, cy + Inches(0.05), cw, Inches(0.35), t,
                 size=12, bold=True, color=WHITE if colors[i % 4] != GOLD else NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, x, cy + Inches(0.45), cw, ch - Inches(0.45), fill=CLOUD, line=LINE)
        add_text(s, x + Inches(0.10), cy + Inches(0.55), cw - Inches(0.2), ch - Inches(0.55),
                 d, size=10, color=INK)
    deal_y = Inches(4.95)
    add_text(s, Inches(0.5), deal_y, Inches(12.5), Inches(0.4),
             "报价与权益（草案）", size=14, bold=True, color=NAVY)
    add_rect(s, Inches(0.5), deal_y + Inches(0.40), Inches(12.5), Emu(20000), fill=GOLD)
    add_rect(s, Inches(0.5), deal_y + Inches(0.55), Inches(8.0), Inches(1.50),
             fill=CLOUD, line=LINE)
    for i, line in enumerate(deal_terms):
        y = deal_y + Inches(0.65) + Inches(0.22) * i
        add_round(s, Inches(0.65), y + Inches(0.05), Inches(0.18), Inches(0.18), "●",
                  fill=BLUE, color=WHITE, size=8)
        add_text(s, Inches(0.90), y, Inches(7.6), Inches(0.22), line,
                 size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(8.65), deal_y + Inches(0.55), Inches(4.35), Inches(1.50),
             fill=NAVY)
    add_text(s, Inches(8.80), deal_y + Inches(0.62), Inches(4.1), Inches(0.32),
             "下一步行动", size=12, bold=True, color=GOLD)
    for i, ln in enumerate(next_steps):
        y = deal_y + Inches(0.95) + Inches(0.32) * i
        add_round(s, Inches(8.80), y + Inches(0.04), Inches(0.7), Inches(0.22),
                  ln[0], fill=GOLD, color=NAVY, size=10, bold=True)
        add_text(s, Inches(9.55), y, Inches(3.4), Inches(0.30), ln[1],
                 size=10, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    return s


# ============================ 主流程 ============================
def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    sw, sh = prs.slide_width, prs.slide_height

    # ============ 1. 封面 ============
    s = new_slide(prs)
    add_rect(s, 0, 0, sw, sh, fill=NAVY)
    add_rect(s, 0, Inches(5.5), sw, Inches(0.06), fill=GOLD)
    deco = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(-2.0), Inches(8.0), Inches(8.0))
    deco.fill.solid(); deco.fill.fore_color.rgb = BLUE
    deco.line.fill.background(); deco.shadow.inherit = False
    add_text(s, Inches(0.8), Inches(0.6), Inches(6), Inches(0.5),
             "SENMA · MODE Hub", size=18, bold=True, color=GOLD)
    add_text(s, Inches(0.8), Inches(1.6), Inches(11), Inches(1.6),
             "森马时尚生态总部园区\n招商方案", size=44, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.8), Inches(11), Inches(0.6),
             "让中国时尚走向世界",
             size=22, color=CLOUD, italic=True)
    add_text(s, Inches(0.8), Inches(4.6), Inches(11), Inches(0.5),
             "5 个 Phase · 8 个任务 · 35 页商务汇报版（与冠松同结构）",
             size=14, color=CLOUD)
    add_round(s, Inches(0.8), Inches(6.2), Inches(2.4), Inches(0.45),
              "v1.0 · 招商策划阶段", fill=GOLD, color=NAVY, size=12, bold=True)
    add_text(s, Inches(3.4), Inches(6.2), Inches(8), Inches(0.45),
             "汇报对象：集团董事会 / 区政府 / 链主品牌",
             size=11, color=CLOUD, anchor=MSO_ANCHOR.MIDDLE)

    # ============ 2. 议程 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="议程",
               page_title="汇报议程", subtitle="约 35 分钟形成完整认知")
    items = [
        ("01", "执行摘要 · 项目概要 · SWOT", "Executive Summary"),
        ("02", "Phase 1 · 策略与定位", "时尚产业链 + 楼栋方案 + 待实测"),
        ("03", "Phase 2 · 招商执行", "链主 5 张一页纸 + 漏斗 + 政府"),
        ("04", "Phase 3 · 品牌与活动", "9 月时装秀 + 12 场活动 + 媒体"),
        ("05", "Phase 4 · 商业条款", "四档 + 服务包 + 报价 + 财务"),
        ("06", "Phase 5 · 落地推进", "12 个月甘特 + 团队 + 风险矩阵"),
        ("07", "投决建议", "三个里程碑承诺 + 5 项授权"),
    ]
    y0 = Inches(1.3)
    for i, (no, title, sub) in enumerate(items):
        y = y0 + Inches(0.75) * i
        add_round(s, Inches(0.8), y, Inches(0.7), Inches(0.55), no,
                  fill=NAVY, color=WHITE, size=18, bold=True)
        add_text(s, Inches(1.7), y, Inches(6), Inches(0.55), title,
                 size=18, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(7.7), y, Inches(5), Inches(0.55), sub,
                 size=13, color=GREY, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(0.8), y + Inches(0.62), Inches(11.7), Emu(15000), fill=LINE)

    # ============ 3. 一页摘要 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="01 · 核心结论",
               page_title="一页看懂：项目战略与三年目标",
               subtitle="SENMA · MODE Hub · 中国时尚生态总部首选地")
    card_w = Inches(2.95); card_h = Inches(1.55); gap = Inches(0.15)
    cards = [
        ("5 万㎡", "1#–5# 综合体（待实测）\n出租净面积约 3.5 万㎡", BLUE),
        ("T 台 ×3", "主秀场 800 人 + SHOWROOM 250\n+ 拍摄棚 ×10", GOLD),
        ("80 间", "直播间 / 跨境运营房\n300 家生态 + 5 链主", PURPLE),
        ("92%", "Y3 入驻率 · 链主 ≥ 3\n设计师 55 + Y3 EBITDA ≈ 1.35 亿", GREEN),
    ]
    x = Inches(0.5)
    for i, (big, sub, color) in enumerate(cards):
        bx = x + (card_w + gap) * i
        add_rect(s, bx, Inches(1.1), card_w, card_h, fill=color)
        add_text(s, bx, Inches(1.18), card_w, Inches(0.85), big,
                 size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, bx, Inches(1.95), card_w, Inches(0.65), sub,
                 size=11, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.5), Inches(2.95), Inches(12.3), Inches(0.4),
             "战略定位 · 三圈层客群 · 五大差异化壁垒",
             size=15, bold=True, color=NAVY)
    add_rect(s, Inches(0.5), Inches(3.40), Inches(12.3), Emu(20000), fill=GOLD)
    block_y = Inches(3.55); block_h = Inches(2.6)

    add_rect(s, Inches(0.5), block_y, Inches(4.0), block_h, fill=CLOUD)
    add_text(s, Inches(0.65), block_y + Inches(0.1), Inches(3.7), Inches(0.4),
             "战略定位", size=14, bold=True, color=NAVY)
    add_text(s, Inches(0.65), block_y + Inches(0.55), Inches(3.7), Inches(2.0),
             "童装 + 设计师 + 时尚电商\n首选总部基地\n\n总部 + 设计 + 电商 + 秀场\n四位一体园区\n\n品牌主张：让中国时尚\n走向世界",
             size=11, color=INK)

    add_rect(s, Inches(4.65), block_y, Inches(4.0), block_h, fill=CLOUD)
    add_text(s, Inches(4.80), block_y + Inches(0.1), Inches(3.7), Inches(0.4),
             "三圈层客群", size=14, bold=True, color=NAVY)
    add_text(s, Inches(4.80), block_y + Inches(0.55), Inches(3.7), Inches(2.0),
             "内圈 · 1# 主楼\n链主品牌总部 · 3–5 家\n\n中圈 · 2#+3#\n50+ 设计师 + 80+ 直播间\n+ 30+ 跨境\n\n外圈 · 4#+5#\nT 台 / 拍摄 / 童装亲子",
             size=11, color=INK)

    add_rect(s, Inches(8.80), block_y, Inches(4.0), block_h, fill=NAVY)
    add_text(s, Inches(8.95), block_y + Inches(0.1), Inches(3.7), Inches(0.4),
             "五大差异化壁垒", size=14, bold=True, color=GOLD)
    add_text(s, Inches(8.95), block_y + Inches(0.55), Inches(3.7), Inches(2.0),
             "① 森马童装生态龙头势能\n\n② 设计师+直播+跨境一站式\n\n③ T 台+SHOWROOM+拍摄棚\n\n④ 童装家庭场景独家\n\n⑤ 政府时尚产业政策包",
             size=11, color=WHITE)

    add_text(s, Inches(0.5), Inches(6.30), Inches(12.3), Inches(0.4),
             "关键举措：5 链主 · 300 生态 · 5 中介 · 1 场 9 月时装秀 · 4 档商业模式",
             size=12, bold=True, color=NAVY)

    # ============ 4. 项目概要（待实测） ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="01 · 项目概要",
               page_title="项目概要 · 5 万方综合体（待实测）",
               subtitle="项目载体由项目方提供建筑设计文件后，将进入 v1.1 实测对齐版")
    # 左：1#-5# 楼栋示意
    add_text(s, Inches(0.5), Inches(1.15), Inches(7.5), Inches(0.4),
             "1#–5# 综合体功能拆分（示意 · 待复核）", size=14, bold=True, color=NAVY)
    fl_x = Inches(0.5); fl_y = Inches(1.55); fl_w = Inches(7.5); fl_h = Inches(0.85)
    blocks = [
        ("1# 主楼 · 时尚总部", "22,000 ㎡ · 9F · 链主总部", NAVY),
        ("2# 设计师楼", "10,000 ㎡ · 6F · 50+ 设计师", BLUE),
        ("3# 电商基地", "8,000 ㎡ · 5F · 80+ 直播间 + 跨境", PURPLE),
        ("4# T 台秀场 + 拍摄棚", "6,000 ㎡ · 主秀 800 人 + 棚 ×10", GOLD),
        ("5# 配套商业 · 童装亲子", "4,000 ㎡ · 集合店 + 母婴 + 咖啡", GREEN),
    ]
    for i, (lbl, desc, c) in enumerate(blocks):
        y = fl_y + (fl_h + Inches(0.05)) * i
        add_rect(s, fl_x, y, Inches(2.5), fl_h, fill=c)
        add_text(s, fl_x, y, Inches(2.5), fl_h, lbl,
                 size=12, bold=True, color=WHITE if c != GOLD else NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, fl_x + Inches(2.55), y, fl_w - Inches(2.55), fl_h, fill=CLOUD, line=LINE)
        add_text(s, fl_x + Inches(2.7), y, fl_w - Inches(2.7), fl_h, desc,
                 size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # 右：关键数据卡
    add_text(s, Inches(8.4), Inches(1.15), Inches(4.7), Inches(0.4),
             "关键数据（招商口径，待实测）", size=14, bold=True, color=NAVY)
    facts = [
        ("总建面",   "约 50,000 ㎡（综合体）", BLUE),
        ("出租净面积", "约 35,000 ㎡", BLUE),
        ("T 台主秀场", "800 人 / SHOWROOM 250", GOLD),
        ("拍摄棚",   "10 个 · 200–800 ㎡", GOLD),
        ("直播间",   "标准 30 ㎡ ×80 + 大型 100 ㎡ ×16", PURPLE),
        ("童装亲子区", "1,500 ㎡（1F+2F）", GREEN),
        ("车位",     "700 个（网红/接送 80）", NAVY),
        ("位置假设", "上海长宁/虹桥 OR 杭州下沙", GREY),
    ]
    for i, (t, d, c) in enumerate(facts):
        y = Inches(1.55) + Inches(0.65) * i
        add_rect(s, Inches(8.4), y, Inches(0.18), Inches(0.55), fill=c)
        add_rect(s, Inches(8.6), y, Inches(4.5), Inches(0.55), fill=CLOUD, line=LINE)
        add_text(s, Inches(8.75), y + Inches(0.05), Inches(1.2), Inches(0.45),
                 t, size=12, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(9.95), y + Inches(0.05), Inches(3.10), Inches(0.45),
                 d, size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # ============ 5. SWOT 总览 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="01 · SWOT",
               page_title="SWOT 总览 · 把森马势能变成差异化报价",
               subtitle="把外部机会与森马生态优势耦合，形成不可复制的护城河")
    quad_w = Inches(6.15); quad_h = Inches(2.85)
    quads = [
        ("优势 S", "Strengths", BLUE,
         ["森马童装生态龙头（巴拉巴拉）",
          "杭州/温州/上海三地协同",
          "直播电商 + 跨境出海双轮",
          "T 台 + 拍摄棚 + 设计师集群独家",
          "童装家庭场景配套（亲子动线）"]),
        ("劣势 W", "Weaknesses", GOLD,
         ["单一项目体量需对标头部园区",
          "设计师品牌生命周期短（<3 年）",
          "时尚电商红利收窄",
          "童装出生率影响",
          "跨境政策不确定性"]),
        ("机会 O", "Opportunities", NAVY,
         ["童装家庭场景升级",
          "国潮 IP + 联名爆发",
          "跨境出海（SHEIN/TEMU）红利",
          "数字化柔性快反技术",
          "海外品牌中国化需要本地总部"]),
        ("威胁 T", "Threats", RED,
         ["杭州艺尚 / 大浪时尚虹吸",
          "直播 GMV 同比下滑",
          "设计师品牌存活率 < 30%",
          "公共预算趋严，补贴收紧",
          "童模 / 直播合规监管收紧"]),
    ]
    positions = [(Inches(0.5), Inches(1.20)), (Inches(6.85), Inches(1.20)),
                 (Inches(0.5), Inches(4.20)), (Inches(6.85), Inches(4.20))]
    for (px, py), (t, en, c, bullets) in zip(positions, quads):
        add_rect(s, px, py, quad_w, quad_h, fill=CLOUD, line=LINE)
        add_rect(s, px, py, quad_w, Inches(0.50), fill=c)
        add_text(s, px + Inches(0.15), py + Inches(0.05), quad_w - Inches(0.3), Inches(0.4),
                 t, size=15, bold=True, color=GOLD if c == NAVY else WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, px + quad_w - Inches(2.0), py + Inches(0.05), Inches(1.8), Inches(0.4),
                 en, size=11, color=WHITE, italic=True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        for i, b in enumerate(bullets):
            yy = py + Inches(0.65) + Inches(0.42) * i
            add_round(s, px + Inches(0.20), yy + Inches(0.08), Inches(0.18), Inches(0.18), "●",
                      fill=c, color=WHITE, size=8)
            add_text(s, px + Inches(0.45), yy, quad_w - Inches(0.6), Inches(0.4), b,
                     size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # ============ 6. Phase 1 扉页 ============
    section_cover(prs, "PHASE 1", "策略与定位",
                  "Strategy & Positioning · 奠基",
                  ["任务 1 · 时尚产业链 + 全国园区对标",
                   "任务 2 · 楼栋功能 + T 台 + 拍摄棚 + 童装亲子"])

    # ============ 7. 时尚产业链 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 1 · 任务 1",
               page_title="时尚产业链六层图谱",
               subtitle="客群匹配的『地图』：从数字基础设施到渠道运营")
    layers = [
        ("L6 渠道与运营", "直播电商 / 私域 / 集合店 / 跨境 / 二手时尚", BLUE),
        ("L5 品牌与 IP", "童装品牌 / 设计师品牌 / 国潮 IP / 联名 / 明星同款", NAVY),
        ("L4 设计与创意", "设计师工作室 / 趋势 / 色彩 / 面料 / 工艺", BLUE),
        ("L3 供应链与制造", "智能工厂 / 柔性快反 / 跨境物流 / 海外仓", NAVY),
        ("L2 材料与面辅料", "功能面料 / 环保再生 / 印染 / 拉链/钮扣", BLUE),
        ("L1 数字基础设施", "ERP / DAM / 3D 设计 / AI 生成 / SaaS / 数据中台", NAVY),
    ]
    bar_x = Inches(0.7); bar_w = Inches(8.5); bar_h = Inches(0.65)
    for i, (lbl, desc, c) in enumerate(layers):
        y = Inches(1.25) + (bar_h + Inches(0.10)) * i
        add_rect(s, bar_x, y, Inches(1.6), bar_h, fill=GOLD)
        add_text(s, bar_x, y, Inches(1.6), bar_h, lbl,
                 size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, bar_x + Inches(1.6), y, bar_w - Inches(1.6), bar_h, fill=c)
        add_text(s, bar_x + Inches(1.7), y, bar_w - Inches(1.7), bar_h, desc,
                 size=11, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    side_x = Inches(9.5); side_w = Inches(3.4)
    side = [("资本 / 孵化", "时尚基金 / 启明 / 高瓴消费"),
            ("时尚教育", "东华大学 · 北服 · 中国美院"),
            ("IP / 版权", "国潮联名 / 博物馆 IP / 明星 IP")]
    for i, (t, d) in enumerate(side):
        y = Inches(1.6) + Inches(1.45) * i
        add_rect(s, side_x, y, side_w, Inches(1.25), fill=CLOUD, line=LINE)
        add_text(s, side_x + Inches(0.15), y + Inches(0.1), side_w - Inches(0.3), Inches(0.4),
                 t, size=14, bold=True, color=NAVY)
        add_text(s, side_x + Inches(0.15), y + Inches(0.5), side_w - Inches(0.3), Inches(0.7),
                 d, size=11, color=INK)
    add_text(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.4),
             "→ 1# 链主品牌 · 2# 设计师/趋势 · 3# 直播+跨境 · 4# T 台/拍摄 · 5# 童装亲子",
             size=12, bold=True, color=BLUE)

    # ============ 8. 全国园区对标 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 1 · 任务 1",
               page_title="全国时尚园区对标矩阵",
               subtitle="结论：时尚总部 + 设计师 + 直播 + T 台四位一体的稀缺组合")
    header = ["园区", "区位", "主导产业", "租金 (元/㎡·天)", "入驻率", "政策亮点", "对我方"]
    rows = [
        ["杭州 · 艺尚小镇", "临平", "时尚总部 / 设计师 / T 台", "4.0–6.5", "~92%", "区+镇双补贴 / 设计师扶持", "离市区远，童装弱"],
        ["上海 · M50 创意园", "普陀", "设计师 / 艺术", "7.0–9.5", "~95%", "文创园区税收减免", "体量小 / 童装弱"],
        ["杭州 · 下沙国际时尚园", "下沙", "时尚电商 / 跨境", "4.5–7.0", "~88%", "跨境补贴 / 电商扶持", "缺顶级品牌总部"],
        ["广州 · 沙河顶时尚谷", "天河", "服装 / 设计师 / 直播", "5.5–8.0", "~85%", "一般", "偏工厂气质"],
        ["深圳 · 大浪时尚小镇", "龙华", "服装总部 / 设计", "5.0–7.5", "~90%", "区级补贴", "设计师密度低"],
        ["上海 · 静安国际时尚中心", "静安", "时装秀 / 高端品牌", "9.0–12.0", "~95%", "国际化资源", "单价高 / 童装弱"],
        ["北京 · 798 / 751 时尚区", "朝阳", "设计师 / 艺术", "8.0–11.0", "~95%", "文创园区", "童装/电商弱"],
    ]
    add_table(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(4.6), header, rows,
              col_widths=[Inches(2.4), Inches(0.8), Inches(1.5), Inches(1.4),
                          Inches(0.9), Inches(2.7), Inches(2.6)])
    add_rect(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.95), fill=NAVY)
    add_text(s, Inches(0.7), Inches(6.05), Inches(12.0), Inches(0.4),
             "我方差异化定价策略", size=13, bold=True, color=GOLD)
    add_text(s, Inches(0.7), Inches(6.40), Inches(12.0), Inches(0.55),
             "基础租金 5.5–9.0 + 流量分发 + 童装亲子协同 + T 台 + 跨境合规一站式 — 等效净价低于核心商圈 30%",
             size=11, color=WHITE)

    # ============ 9. 楼栋平面 + 面积表 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 1 · 任务 2",
               page_title="1#–5# 楼栋功能拆分 + 时尚资产",
               subtitle="链主总部 / 设计师 / 电商 / T 台 + 拍摄棚 / 童装亲子（待实测）")
    base_x = Inches(0.5); base_y = Inches(1.2)
    plot_w = Inches(7.0); plot_h = Inches(5.4)
    add_rect(s, base_x, base_y, plot_w, plot_h, fill=CLOUD, line=LINE)
    add_text(s, base_x, base_y + Inches(0.05), plot_w, Inches(0.3),
             "园区平面示意（北 ↑）", size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    bx, by = base_x + Inches(0.4), base_y + Inches(0.5)
    bw, bh = Inches(1.5), Inches(1.4)
    add_rect(s, bx, by, bw, bh, fill=NAVY)
    add_text(s, bx, by, bw, bh, "1# 主楼\n时尚总部\n22,000 ㎡",
             size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    ax = base_x + plot_w - bw - Inches(0.4)
    add_rect(s, ax, by, bw, bh, fill=GOLD)
    add_text(s, ax, by, bw, bh, "2# 设计师楼\n10,000 ㎡",
             size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cx2 = bx + bw + Inches(0.3); cw2 = ax - cx2 - Inches(0.3)
    add_rect(s, cx2, by, cw2, bh, fill=WHITE, line=LINE)
    add_text(s, cx2, by, cw2, bh, "中央广场 + 红毯入口\n林荫水景",
             size=10, color=GREY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cy2 = by + bh + Inches(0.2)
    add_rect(s, bx, cy2, bw, bh, fill=PURPLE)
    add_text(s, bx, cy2, bw, bh, "3# 电商基地\n8,000 ㎡",
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, ax, cy2, bw, bh, fill=BLUE)
    add_text(s, ax, cy2, bw, bh, "5# 配套\n童装亲子\n4,000 ㎡",
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, cx2, cy2 + Inches(0.5), cw2, Inches(0.4), fill=CLOUD, line=LINE)
    add_text(s, cx2, cy2 + Inches(0.5), cw2, Inches(0.4), "连廊 / 食堂",
             size=10, color=GREY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    ty = cy2 + bh + Inches(0.25)
    test_w = bw + cw2 + Inches(0.3)
    add_rect(s, bx, ty, test_w, Inches(1.6), fill=GREEN)
    add_text(s, bx, ty, test_w, Inches(1.6),
             "4# T 台秀场 + 拍摄棚  6,000 ㎡\n主秀 800 人 + SHOWROOM 250 + 拍摄棚 ×10",
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, ax, ty, bw, Inches(1.6), fill=NAVY)
    add_text(s, ax, ty, bw, Inches(1.6), "停车 + 库\n网红车位 80",
             size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    header = ["区域", "面积 (㎡)", "性质"]
    rows = [
        ["1# 时尚总部", "22,000", "出租"],
        ["2# 设计师楼", "10,000", "出租 + 工位"],
        ["3# 电商基地", "8,000", "出租 + 直播位"],
        ["4# T 台 + 拍摄", "6,000", "公共/会员"],
        ["5# 配套商业", "4,000", "部分自留"],
        ["中央广场/连廊", "3,000", "公共"],
        ["食堂/健身/母婴", "2,000", "公共"],
        ["地下停车/库", "12,000", "商业"],
        ["合计建筑面积", "约 67,000", "—"],
        ["可出租净面积", "约 35,000", "—"],
    ]
    add_table(s, Inches(7.8), Inches(1.2), Inches(5.0), Inches(5.4),
              header, rows,
              col_widths=[Inches(2.0), Inches(1.4), Inches(1.6)])

    # ============ 10. 1# 主楼 链主总部深度 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 1 · 任务 2",
               page_title="1# 主楼 · 链主品牌总部（『时尚灯塔』）",
               subtitle="9F · 22,000 ㎡ · 起始 5.5–7.0 元/㎡·天 · 5–8 年长租")
    add_text(s, Inches(0.5), Inches(1.15), Inches(5), Inches(0.4),
             "纵剖示意（9F · 22,000 ㎡）", size=14, bold=True, color=NAVY)
    fl_x = Inches(0.5); fl_y = Inches(1.55); fl_w = Inches(5.5); fl_h = Inches(0.50)
    floors = [
        ("9F", "屋顶花园 + CEO 接待 / 露台", GOLD),
        ("8F", "战略 / 投资 / 高管", NAVY),
        ("7F", "品牌总部办公（链主 A）", BLUE),
        ("6F", "品牌总部办公（链主 A）", BLUE),
        ("5F", "品牌总部办公（链主 B）", PURPLE),
        ("4F", "品牌总部办公（链主 B）", PURPLE),
        ("3F", "童装研发 / 设计中心", GREEN),
        ("2F", "客户接待 / 路演 / 集合店", GOLD),
        ("1F", "礼宾大堂 + 集合店主入口", NAVY),
        ("B1", "停车 / 设备 / 库", GREY),
    ]
    for i, (lvl, desc, c) in enumerate(floors):
        y = fl_y + (fl_h + Inches(0.04)) * i
        add_rect(s, fl_x, y, Inches(0.85), fl_h, fill=c)
        add_text(s, fl_x, y, Inches(0.85), fl_h, lvl,
                 size=11, bold=True, color=WHITE if c not in (GOLD, GREY) else NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, fl_x + Inches(0.9), y, fl_w - Inches(0.9), fl_h, fill=CLOUD, line=LINE)
        add_text(s, fl_x + Inches(1.0), y, fl_w - Inches(1.0), fl_h, desc,
                 size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(7.0), Inches(1.15), Inches(6), Inches(0.4),
             "核心卖点", size=14, bold=True, color=NAVY)
    sells = [
        ("楼层冠名权", "「[品牌] · MODE Hub」 5 年", GOLD),
        ("时尚集合店主入口", "1F 红毯 + 5.5 m 挑空大堂", BLUE),
        ("9F 屋顶花园", "T 台彩排 / ESG 媒体大片", NAVY),
        ("童装家庭场景", "1F 童装亲子区直接对接", GREEN),
        ("配套联动", "T 台 / 拍摄棚 / 设计师工作室", PURPLE),
        ("机电预留", "≥ 200 W/㎡ · 4.5 m 净高", BLUE),
    ]
    for i, (t, d, c) in enumerate(sells):
        y = Inches(1.55) + Inches(0.55) * i
        add_rect(s, Inches(7.0), y, Inches(0.18), Inches(0.45), fill=c)
        add_rect(s, Inches(7.2), y, Inches(5.7), Inches(0.45), fill=CLOUD, line=LINE)
        add_text(s, Inches(7.35), y + Inches(0.02), Inches(2.3), Inches(0.41),
                 t, size=11, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(9.65), y + Inches(0.02), Inches(3.2), Inches(0.41),
                 d, size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    add_rect(s, Inches(7.0), Inches(4.95), Inches(5.9), Inches(1.95), fill=NAVY)
    add_text(s, Inches(7.15), Inches(5.0), Inches(5.6), Inches(0.4),
             "签约组合（草案）", size=13, bold=True, color=GOLD)
    deal = ("起始租金 5.5–7.0 元/㎡·天 · 5–8 年长租\n"
            "免租 12–18 个月 · 装补 800–1,200 元/㎡（封顶 400 万）\n"
            "区级税收留成 80% 三年返 / 50% 后两年返\n"
            "T 台主秀场 + 拍摄棚终身免费会员\n"
            "人才公寓 100 套 + 落户绿通 30 个/年\n"
            "9 月时装秀主秀场冠名 + 媒体首发权")
    add_text(s, Inches(7.15), Inches(5.35), Inches(5.6), Inches(1.55),
             deal, size=10, color=WHITE)

    # ============ 11. 2#+3# 设计师 + 电商 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 1 · 任务 2",
               page_title="2# 设计师楼 + 3# 电商基地",
               subtitle="设计师堆栈层 · 直播 + 跨境双发动机")
    cards = [
        ("2# 设计师楼\n10,000 ㎡ · 6F", "7.0–8.5 元/㎡·天 + 工位",
         ["50+ 独立设计师工作室",
          "面料图书馆 5,000+ 样品",
          "3D 打样 + 数码印花 + 小批量打版",
          "设计师集合店 + 联名快闪",
          "签约 3+3+1 年 · 免租 3–6 月 · 装补 300–500/㎡"], BLUE),
        ("3# 电商基地\n8,000 ㎡ · 5F", "8.0–9.0 元/㎡·天 + 直播位",
         ["80 个标准直播间（30 ㎡）+ 16 个大型（100 ㎡）",
          "跨境运营房 200–500 ㎡ × 多家",
          "网红工位 500 元/月 + 共享拍摄",
          "直播一键开播 + 跨境合规中心",
          "签约月付制 · 服务佣金 GMV 5%–10%"], PURPLE),
        ("4# T 台 + 拍摄棚\n6,000 ㎡", "会员 / 按场次 / 月卡",
         ["主秀场 500–800 人 T 台",
          "SHOWROOM 200–250 人 ×2",
          "拍摄棚 ×10（800/400/200 ㎡）",
          "化妆 / 后台 / 道具 / 直播间",
          "链主免费 / 设计师 5k–30k/场"], GOLD),
    ]
    cw = Inches(4.10); cy = Inches(1.20); ch = Inches(5.5); gx = Inches(0.10)
    for i, (t, sub, bullets, c) in enumerate(cards):
        x = Inches(0.5) + (cw + gx) * i
        add_rect(s, x, cy, cw, ch, fill=CLOUD, line=LINE)
        add_rect(s, x, cy, cw, Inches(1.05), fill=c)
        add_text(s, x, cy + Inches(0.10), cw, Inches(0.5), t,
                 size=18, bold=True, color=NAVY if c == GOLD else WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, cy + Inches(0.55), cw, Inches(0.4), sub,
                 size=12, color=NAVY if c == GOLD else WHITE, align=PP_ALIGN.CENTER)
        for j, b in enumerate(bullets):
            yy = cy + Inches(1.30) + Inches(0.85) * j
            add_round(s, x + Inches(0.20), yy + Inches(0.18), Inches(0.30), Inches(0.30),
                      str(j + 1), fill=c, color=WHITE if c != GOLD else NAVY,
                      size=12, bold=True)
            add_text(s, x + Inches(0.55), yy, cw - Inches(0.7), Inches(0.85),
                     b, size=11, color=INK)

    # ============ 12. 5# 童装亲子 + 时尚资产 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 1 · 任务 2",
               page_title="5# 配套商业 · 童装亲子区 + 时尚资产成绩单",
               subtitle="独家壁垒：森马童装生态 + 设计师 + T 台 + 跨境一站式")
    add_text(s, Inches(0.5), Inches(1.15), Inches(6), Inches(0.4),
             "5# 配套商业 · 童装亲子区 / 集合店 / 咖啡", size=14, bold=True, color=NAVY)
    e_floors = [
        ("2F", "时尚集合店 + 设计师品牌精品", GOLD),
        ("1F", "童装亲子区 + 巴拉巴拉旗舰\n+ 早教 + 摄影 + 体验店", GREEN),
        ("外延", "5 间精品咖啡 + Lounge\n+ 海外品牌方驻华办事处", BLUE),
    ]
    for i, (lvl, d, c) in enumerate(e_floors):
        y = Inches(1.55) + Inches(1.40) * i
        add_rect(s, Inches(0.5), y, Inches(1.0), Inches(1.30), fill=c)
        add_text(s, Inches(0.5), y, Inches(1.0), Inches(1.30), lvl,
                 size=14, bold=True, color=NAVY if c == GOLD else WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(1.55), y, Inches(4.95), Inches(1.30), fill=CLOUD, line=LINE)
        add_text(s, Inches(1.70), y + Inches(0.05), Inches(4.7), Inches(1.20),
                 d, size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(6.85), Inches(1.15), Inches(6.2), Inches(0.4),
             "时尚资产 · 一站式成绩单", size=14, bold=True, color=NAVY)
    scenes = [
        ("T 台主秀场",     "800 人 + 设计师奖典礼", GOLD),
        ("SHOWROOM ×2",    "250 人 / 多功能", BLUE),
        ("拍摄棚 ×10",     "800/400/200 ㎡ + 化妆 + 道具", PURPLE),
        ("直播间 ×80",     "标准 + 大型 + 跨境运营房", BLUE),
        ("面料图书馆",     "5,000+ 样品 / 季度更新", GREEN),
        ("3D 打样",        "数码印花 + 小批量打版", GREEN),
        ("童装家庭场景",   "1,500 ㎡ · 早教 + 体验", GOLD),
    ]
    for i, (t, d, c) in enumerate(scenes):
        y = Inches(1.55) + Inches(0.65) * i
        add_rect(s, Inches(6.85), y, Inches(0.18), Inches(0.55), fill=c)
        add_rect(s, Inches(7.05), y, Inches(6.0), Inches(0.55), fill=CLOUD, line=LINE)
        add_text(s, Inches(7.20), y + Inches(0.05), Inches(2.6), Inches(0.45),
                 t, size=11, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(9.85), y + Inches(0.05), Inches(3.1), Inches(0.45),
                 d, size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(6.85), Inches(6.20), Inches(6.2), Inches(0.55), fill=NAVY)
    add_text(s, Inches(7.0), Inches(6.25), Inches(6.0), Inches(0.45),
             "→ 把『总部+设计+电商+秀场』四位一体的稀缺资产做成林立护城河",
             size=11, bold=True, color=GOLD, anchor=MSO_ANCHOR.MIDDLE)

    # ============ 13. Phase 2 扉页 ============
    section_cover(prs, "PHASE 2", "招商执行（核心）",
                  "Leasing Execution · Anchor + Ecosystem + Government",
                  ["任务 3 · 链主 TOP5 攻坚 + 5 张一页纸",
                   "任务 4 · 300 家生态招商漏斗",
                   "任务 5 · 政府关系与政策包"])

    # ============ 14. 链主 TOP5 作战图 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 2 · 任务 3",
               page_title="链主 TOP5 攻坚作战图",
               subtitle="项目总监 + BD 总监 双人出动 · 周一例会 / 季度董事会复盘")
    header = ["#", "链主", "落位", "决策人 / 关键人", "当前阶段", "下一动作", "截止", "状态"]
    rows = [
        ["1", "江南布衣 JNBY", "1# 8F 整层 ~2,500 ㎡", "李琳 / 吴健", "T0 名片建联", "中国服装协会引荐 + 走访", "T+30d", "推进中"],
        ["2", "之禾 ICICLE", "1# 9F 整层 ~2,400 ㎡", "叶寿增", "已发邀约函", "时装秀主秀场邀请", "T+45d", "顺利"],
        ["3", "安踏儿童 FILA Kids", "1# 8F 整层 ~1,800 ㎡", "丁世忠 / 林翔华", "公司层初接触", "园区参访 + 巴拉巴拉协同", "T+60d", "推进中"],
        ["4", "Carter's / GAP Kids", "1# 7F 部分 ~1,500 ㎡", "中国 GM", "高层已建联", "外资协会引荐 / 跨境合规", "T+30d", "顺利"],
        ["5", "SHEIN 设计中心", "2# + 3# 共 ~4,500 ㎡", "中国设计 VP", "已建联，等会面", "30 分钟会面", "T+45d", "推进中"],
    ]
    add_table(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.6),
              header, rows,
              col_widths=[Inches(0.4), Inches(1.7), Inches(2.1), Inches(1.8),
                          Inches(1.6), Inches(2.6), Inches(1.0), Inches(1.1)])
    add_text(s, Inches(0.5), Inches(4.0), Inches(12), Inches(0.4),
             "每家「一页纸」定制提案 · 四要点结构（详见接下来 5 页）",
             size=14, bold=True, color=NAVY)
    add_rect(s, Inches(0.5), Inches(4.42), Inches(12.3), Emu(20000), fill=GOLD)
    quad_y = Inches(4.6); quad_h2 = Inches(2.3)
    quads2 = [
        ("决策人 / 关键人", "30 秒电梯演讲 → CEO/创始人；BD 把行业协会接口排进议程", BLUE),
        ("正缺什么", "上海 / 长三角双总部 · 时装秀资源 · 跨境合规 · 童装协同", NAVY),
        ("独家供给", "楼层冠名 · T 台 + 拍摄棚 · 童装亲子 · 设计师集合店", GOLD),
        ("报价权益", "起始租 5.5–7.0 · 免租 12–18m · 装补 800–1,200/㎡ · 政策返还 80%", GREEN),
    ]
    qw = Inches(2.95); gx = Inches(0.15)
    for i, (t, d, c) in enumerate(quads2):
        x = Inches(0.5) + (qw + gx) * i
        add_rect(s, x, quad_y, qw, quad_h2, fill=c)
        add_text(s, x + Inches(0.15), quad_y + Inches(0.15), qw - Inches(0.3), Inches(0.5),
                 t, size=14, bold=True, color=NAVY if c == GOLD else WHITE)
        add_text(s, x + Inches(0.15), quad_y + Inches(0.7), qw - Inches(0.3), quad_h2 - Inches(0.85),
                 d, size=11, color=NAVY if c == GOLD else WHITE)

    # ============ 15-19. 链主 5 张一页纸 ============
    anchor_one_pager(prs, idx=1, label="JNBY", brand="江南布衣 JNBY",
        title="上海长三角双总部 + 设计师生态共建",
        why_anchor="中国设计师品牌龙头，杭州总部已稳。\n需上海/长三角第二总部辐射华东 + 海外。\n与森马生态高度互补（成人时尚 + 童装）。",
        what_lacks="上海设计师人才 + 海外资源\n童装板块协同（蓬马童装）\n时装秀 / SHOWROOM 资源\n跨境出海生态接口",
        unique_offer=[
            ("8F 整层 + 楼层冠名", "「JNBY · MODE Hub」 5 年"),
            ("设计师集群联动", "50+ 设计师 + WGSN/POP"),
            ("T 台主秀场", "800 人 + 拍摄棚 ×10"),
            ("巴拉巴拉数据闭环", "童装协同 + 跨境出海"),
        ],
        deal_terms=[
            "1# 8F 整层 ~2,500 ㎡ · 签约 6 年",
            "起始租金 6.0 元/㎡·天 · 3 年一调",
            "免租 14 个月 · 装补 1,000 元/㎡ (封顶 250 万)",
            "T 台主秀首场冠名 + 拍摄棚 ≥ 50 场次/年免费",
            "区级税收 80% 三年返 + 人才公寓 100 套",
        ],
        next_steps=[
            ("T+15d", "中国服装协会引荐 + 李琳走访"),
            ("T+30d", "李琳 + 集团董事长互访"),
            ("T+90d", "意向条款书签署"),
        ])

    anchor_one_pager(prs, idx=2, label="之禾", brand="之禾 ICICLE",
        title="海外大本营延伸 + 设计师 + 高端生活方式",
        why_anchor="中国高端时装代表，海外（巴黎/米兰）已稳。\n需要在国内核心圈层做「高端生活方式总部」叙事。",
        what_lacks="高端品牌总部展示空间\n设计师生态协同\n海外业务接待\n跨境出海一站式",
        unique_offer=[
            ("9F 整层 + 屋顶花园", "高端品牌门面"),
            ("1F 集合店主入口", "高端时装陈列"),
            ("T 台主秀场冠名", "年度时装秀主旨发布"),
            ("跨境合规一站式", "SHEIN/TEMU/TikTok"),
        ],
        deal_terms=[
            "1# 9F 整层 ~2,400 ㎡ · 签约 5 年",
            "起始租金 6.5 元/㎡·天",
            "免租 12 个月 · 装补 1,000 元/㎡",
            "屋顶花园专属 + T 台 ≥ 30 场次/年免费",
            "海外业务接待 + 酒店式服务",
        ],
        next_steps=[
            ("T+10d", "邀请叶寿增参加时装秀主秀场"),
            ("T+30d", "选址踏勘 + Term Sheet"),
            ("T+90d", "签约"),
        ])

    anchor_one_pager(prs, idx=3, label="安踏儿童", brand="安踏儿童 FILA Kids",
        title="中国童装上海设计中心 + 巴拉巴拉协同",
        why_anchor="中国童装第二阵营龙头，急需在上海建立设计 + 营销双中心。\n与森马巴拉巴拉形成「中国童装双雄」叙事。",
        what_lacks="上海童装设计中心\n童装家庭场景体验\n设计师 + 模特资源\n直播电商基地",
        unique_offer=[
            ("1# 8F 整层 + 楼层冠名", "「FILA Kids · MODE Hub」"),
            ("童装亲子区协同", "1F 早教 + 摄影 + 体验"),
            ("拍摄棚 ×3 长期使用", "设计师 + 童模库"),
            ("80 个直播间", "跨境出海一站式"),
        ],
        deal_terms=[
            "1# 8F 部分 ~1,800 ㎡ · 签约 5 年",
            "起始租金 6.0 元/㎡·天",
            "免租 12 个月 · 装补 800 元/㎡",
            "童装亲子区 1F 协同合作",
            "区级税收 80% 三年返",
        ],
        next_steps=[
            ("T+15d", "林翔华参访园区"),
            ("T+45d", "与巴拉巴拉协同会"),
            ("T+90d", "意向书签署"),
        ])

    anchor_one_pager(prs, idx=4, label="外资童装", brand="Carter's / GAP Kids",
        title="外资童装中国总部 + 跨境合规接口",
        why_anchor="美国童装第一品牌（Carter's）+ 大众童装代表（GAP Kids）。\n需要在中国设立『区域总部 + 设计 + 跨境』复合中心。",
        what_lacks="中国区跨境合规接口\n童装生态协同\n设计本土化\n外资政策对接",
        unique_offer=[
            ("跨境合规中心", "关务+税务+知识产权"),
            ("童装亲子区", "巴拉巴拉数据 + 早教"),
            ("50+ 设计师集群", "趋势/面料图书馆"),
            ("外资协会驻点", "时尚律所 + 政府绿通"),
        ],
        deal_terms=[
            "1# 7F 部分 ~1,500 ㎡ · 签约 5 年",
            "起始租金 6.5 元/㎡·天（外资品牌价）",
            "免租 12 个月 · 装补 600 元/㎡",
            "跨境合规中心 + 律所对接 + 政府绿通",
            "海外品牌方驻华办事处协同",
        ],
        next_steps=[
            ("T+15d", "联合外资协会拜访"),
            ("T+45d", "参访园区与跨境合规中心"),
            ("T+90d", "意向书签署"),
        ])

    anchor_one_pager(prs, idx=5, label="SHEIN", brand="SHEIN 设计中心",
        title="上海设计师生态枢纽 + 拍摄供应链",
        why_anchor="中国跨境时尚电商第一名，全球 GMV 超 4,000 亿。\n急需在上海建立设计师 + 拍摄 + 跨境运营复合中心。",
        what_lacks="上海设计师集群（快反）\n拍摄棚 + 模特\n跨境出海与合规\n网红与直播",
        unique_offer=[
            ("2# 5F 设计 + 3F 跨境", "共 4,500 ㎡"),
            ("拍摄棚 ×10 + 模特库", "≥ 100 场次/年免费"),
            ("跨境合规中心", "TikTok/Shopify 一站式"),
            ("80+ 直播间", "MCN 矩阵协同"),
        ],
        deal_terms=[
            "2# 5F 设计 + 3# 跨境 共 4,500 ㎡ · 签约 5 年",
            "起始租金 5.8 元/㎡·天（链主特惠）",
            "免租 14 个月 · 装补 1,000 元/㎡",
            "拍摄棚 ≥ 100 场次/年免费",
            "设计师协同合作 + 模特库",
        ],
        next_steps=[
            ("T+10d", "SHEIN 设计 VP 30 分钟会面"),
            ("T+45d", "Term Sheet 草案"),
            ("T+120d", "签约"),
        ])

    # ============ 20. 300 家生态漏斗 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 2 · 任务 4",
               page_title="生态企业招商漏斗 · 300 家库",
               subtitle="5 家中介 · 9 类来源 · 端到端转化 17% · 80 家年签约目标")
    funnel = [
        ("L1 线索 (Lead)", 480, 12.0, BLUE),
        ("L2 触达 (Reached)", 384, 9.7, NAVY),
        ("L3 意向 (Intent)", 192, 4.8, BLUE),
        ("L4 谈判 (Negotiation)", 115, 2.9, NAVY),
        ("L5 签约 (Signed)", 80, 2.0, GOLD),
        ("L6 入驻 (Move-in)", 76, 1.9, GREEN),
    ]
    fy = Inches(1.25)
    for i, (lbl, n, w_inch, c) in enumerate(funnel):
        y = fy + Inches(0.65) * i
        add_rect(s, Inches(0.5), y, Inches(3.2), Inches(0.55), fill=CLOUD, line=LINE)
        add_text(s, Inches(0.55), y, Inches(3.1), Inches(0.55), lbl,
                 size=11, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(3.8), y, Inches(w_inch * 0.55), Inches(0.55), fill=c)
        add_text(s, Inches(3.8), y, Inches(w_inch * 0.55) + Inches(1.2), Inches(0.55),
                 f"  {n} 家", size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.5), Inches(5.5), Inches(10), Inches(0.4),
             "理论端到端转化（线索 → 签约）： 80% × 50% × 60% × 70% ≈ 17%",
             size=12, bold=True, color=NAVY)
    add_text(s, Inches(0.5), Inches(5.92), Inches(10), Inches(0.4),
             "80 家签约 ⇐ 480 有效线索 ⇐ 300 库 + 自拓溢出",
             size=12, color=BLUE)
    header = ["中介", "侧重", "佣金档", "年度任务"]
    rows = [
        ["戴德梁行 C&W", "跨国/外资童装", "首月 100% / 长租 120%", "≥ 60 家"],
        ["仲量联行 JLL", "时尚品牌总部", "同上", "≥ 60 家"],
        ["高力 Colliers", "设计师 / 中型时尚", "同上", "≥ 50 家"],
        ["世邦 CBRE", "链主 / 整层", "链主最高 150%", "≥ 40 家"],
        ["本地精品行", "设计师 / MCN / 直播", "首月 80%", "≥ 90 家"],
    ]
    add_text(s, Inches(8.2), Inches(1.0), Inches(5), Inches(0.35),
             "5 家中介渠道 · 非独家 + 30 天首报", size=13, bold=True, color=NAVY)
    add_table(s, Inches(8.2), Inches(1.4), Inches(4.8), Inches(3.1),
              header, rows,
              col_widths=[Inches(1.6), Inches(1.4), Inches(1.0), Inches(0.8)])
    add_text(s, Inches(8.2), Inches(4.7), Inches(5), Inches(0.35),
             "300 家库 · 产业链层级配比", size=13, bold=True, color=NAVY)
    bars = [("L4 设计师", 80), ("L6 直播 MCN", 40), ("L5 设计师品牌", 30), ("L6 跨境", 30),
            ("L5 童装品牌", 20), ("L3 供应链", 20), ("L4 趋势/面料", 20),
            ("L1 数字", 15), ("时尚服务/教育", 30), ("法律/金融", 15)]
    by0 = Inches(5.05)
    for i, (lbl, n) in enumerate(bars):
        y = by0 + Inches(0.20) * i
        add_text(s, Inches(8.2), y, Inches(1.4), Inches(0.18), lbl, size=9, color=INK,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(9.7), y + Inches(0.04), Inches(n / 80.0 * 2.2), Inches(0.12), fill=BLUE)
        add_text(s, Inches(12.0), y, Inches(0.8), Inches(0.18), str(n), size=9, color=NAVY,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)

    # ============ 21. 政府关系 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 2 · 任务 5",
               page_title="政府关系对接 · 路径地图与政策包",
               subtitle="区/市双层 · 7 步首轮汇报 · 形成「一企一策」")
    add_text(s, Inches(0.5), Inches(1.15), Inches(6), Inches(0.4),
             "对接路径地图（市/区双层）", size=14, bold=True, color=NAVY)
    add_rect(s, Inches(0.5), Inches(1.55), Inches(6.2), Inches(5.3), fill=CLOUD, line=LINE)
    nodes = [
        ("市委市政府", Inches(2.5), Inches(1.7), Inches(2.0), Inches(0.45), NAVY, WHITE),
        ("市经信委 (时尚)", Inches(0.7), Inches(2.5), Inches(1.6), Inches(0.45), BLUE, WHITE),
        ("市商务委 (跨境)", Inches(2.5), Inches(2.5), Inches(2.0), Inches(0.45), BLUE, WHITE),
        ("市文旅 / 网信", Inches(4.7), Inches(2.5), Inches(1.8), Inches(0.45), BLUE, WHITE),
        ("区委区政府", Inches(2.3), Inches(3.5), Inches(2.4), Inches(0.5), GOLD, NAVY),
        ("区投促办", Inches(0.7), Inches(4.4), Inches(1.6), Inches(0.45), NAVY, WHITE),
        ("区商务/文旅", Inches(2.5), Inches(4.4), Inches(2.0), Inches(0.45), NAVY, WHITE),
        ("区财政/税务", Inches(4.7), Inches(4.4), Inches(1.8), Inches(0.45), NAVY, WHITE),
        ("区市监 (品牌/直播)", Inches(0.7), Inches(5.1), Inches(1.6), Inches(0.45), BLUE, WHITE),
        ("区人社/房管", Inches(2.5), Inches(5.1), Inches(2.0), Inches(0.45), BLUE, WHITE),
        ("区国资/街道", Inches(4.7), Inches(5.1), Inches(1.8), Inches(0.45), BLUE, WHITE),
    ]
    for txt, x, y, w, h, fc, tc in nodes:
        add_round(s, x, y, w, h, txt, fill=fc, color=tc, size=10, bold=True)
    add_text(s, Inches(0.7), Inches(5.9), Inches(6), Inches(0.4),
             "7 步首轮汇报 · 30 天完成全覆盖", size=12, bold=True, color=NAVY)
    add_text(s, Inches(0.7), Inches(6.25), Inches(6), Inches(0.6),
             "区投促办 → 副区长 → 区委书记/区长 → 市经信委 → 市商务委 → 市文旅/网信 → 区四套班子专题会",
             size=10, color=INK)
    add_text(s, Inches(7.0), Inches(1.15), Inches(6), Inches(0.4),
             "政策包 · 六维适配清单", size=14, bold=True, color=NAVY)
    pol = [
        ("时尚产业专项", "园区认定 · 设计师孵化 · 童装升级 · 时装周分会场", BLUE),
        ("跨境与出海", "跨境综试 · 海外仓 · 退税 · 巴黎/米兰参展补贴", GOLD),
        ("直播与合规", "直播基地认定 · 内容审核绿通 · MCN 培训", PURPLE),
        ("财税返还", "区级留成 80% 三年返 / 50% 后两年返", BLUE),
        ("人才 / 公寓", "链主 100 套 + 设计师 50 套 + 主播 50 套 + 落户 80/年", GREEN),
        ("文化与传播", "时装周 · 时尚电影 · 国潮联名 · 童装家庭", NAVY),
    ]
    for i, (t, d, c) in enumerate(pol):
        y = Inches(1.55) + Inches(0.85) * i
        add_rect(s, Inches(7.0), y, Inches(0.2), Inches(0.75), fill=c)
        add_rect(s, Inches(7.2), y, Inches(5.7), Inches(0.75), fill=CLOUD, line=LINE)
        add_text(s, Inches(7.35), y + Inches(0.05), Inches(5.3), Inches(0.32),
                 t, size=12, bold=True, color=NAVY)
        add_text(s, Inches(7.35), y + Inches(0.36), Inches(5.3), Inches(0.4),
                 d, size=10, color=INK)

    # ============ 22. Phase 3 扉页 ============
    section_cover(prs, "PHASE 3", "品牌与活动",
                  "Brand & Events · Make Noise, Build Trust",
                  ["任务 6 · 9 月年度时装秀 + 12 场活动 + 媒体清单"])

    # ============ 23. 9 月时装秀 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 3 · 任务 6",
               page_title="9 月旗舰：MODE Hub 年度时装秀",
               subtitle="3 天 · 主秀 800 人 ×3 场 + SHOWROOM 250 人 ×6 场 + 在线 100 万")
    agenda = [
        ("Day 1 13:00", "媒体日 · 园区参访 + 设计师工作室开放"),
        ("Day 1 18:00", "媒体晚宴 + 内容预热"),
        ("Day 2 10:00", "开幕式 + 区政府致辞 + 主旨发布"),
        ("Day 2 14:00", "链主品牌主秀 ×3"),
        ("Day 2 19:30", "重大签约仪式 + 招待晚宴"),
        ("Day 3 11:00", "独立设计师集合秀 ×6"),
        ("Day 3 19:00", "设计师奖典礼 + 颁奖"),
    ]
    add_text(s, Inches(0.5), Inches(1.2), Inches(6), Inches(0.4),
             "三天议程", size=14, bold=True, color=NAVY)
    for i, (t, d) in enumerate(agenda):
        y = Inches(1.65) + Inches(0.55) * i
        add_rect(s, Inches(0.5), y, Inches(1.8), Inches(0.45), fill=NAVY)
        add_text(s, Inches(0.5), y, Inches(1.8), Inches(0.45), t,
                 size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(2.35), y, Inches(4.30), Inches(0.45), fill=CLOUD, line=LINE)
        add_text(s, Inches(2.45), y, Inches(4.20), Inches(0.45), d,
                 size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(7.0), Inches(1.2), Inches(6), Inches(0.4),
             "受邀人结构（累计 ~3,000 人次）", size=14, bold=True, color=NAVY)
    invitees = [("政府领导", 50), ("链主 CEO", 80), ("设计师", 300), ("买手/经销商", 400),
                ("媒体/KOL", 300), ("投资机构", 100), ("客户/会员", 1500), ("内部", 270)]
    for i, (lbl, n) in enumerate(invitees):
        y = Inches(1.65) + Inches(0.32) * i
        add_text(s, Inches(7.0), y, Inches(1.6), Inches(0.28), lbl,
                 size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(8.7), y + Inches(0.05), Inches(n / 1500.0 * 3.0), Inches(0.18), fill=GOLD)
        add_text(s, Inches(11.8), y, Inches(0.8), Inches(0.28), str(n),
                 size=10, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)

    add_text(s, Inches(7.0), Inches(4.30), Inches(6), Inches(0.4),
             "倒推时间表（T = 时装秀 Day 1）", size=14, bold=True, color=NAVY)
    timeline = ["T-120d 立项", "T-90d 视觉/秀册 + 邀请函", "T-60d 签约项目锁定 + 设计师终选",
                "T-45d 彩排 1", "T-30d 内容定稿", "T-21d 彩排 2 + 童模合规",
                "T-14d 彩排 3", "T-7d 彩排 4", "T-1d 全员彩排", "T 日 执行 + T+7d 复盘"]
    for i, t in enumerate(timeline):
        y = Inches(4.7) + Inches(0.22) * i
        add_round(s, Inches(7.0), y, Inches(0.3), Inches(0.18), "●",
                  fill=BLUE, color=WHITE, size=8, bold=True)
        add_text(s, Inches(7.4), y - Inches(0.01), Inches(5.5), Inches(0.22),
                 t, size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # ============ 24. 年度活动 + 媒体 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 3 · 任务 6",
               page_title="年度 12 场活动日历 + 媒体四层矩阵",
               subtitle="1 时装秀 + 4 季度沙龙 + 4 月度生态 + 1 跨境峰会 + 2 童装嘉年华")
    header = ["#", "月份", "活动", "形式 / 规模", "主题方向"]
    rows = [
        ["1", "2026-05", "园区品牌发布会（启动）", "现场 200 人", "项目首秀 / 链主官宣"],
        ["2", "2026-06", "童装家庭嘉年华（六一）", "公众 5,000 人", "童装 + 早教 + 亲子"],
        ["3", "2026-07", "直播电商沙龙", "闭门 60 人", "GMV / 跨境 / 合规"],
        ["4", "2026-08", "设计师集合开放日", "半日 80 人", "50 个工作室开放"],
        ["5", "2026-09", "MODE Hub 年度时装秀", "3,000 人次", "链主签约 / 设计师奖"],
        ["6", "2026-10", "跨境出海峰会（国际）", "现场 250 人", "SHEIN/TEMU/TikTok"],
        ["7", "2026-11", "时尚投融资 Demo Day", "半日 80 人", "设计师 / MCN"],
        ["8", "2026-12", "年终招商盘点 + 入驻晚宴", "晚宴 200 人", "全年成果"],
        ["9", "2027-01", "春节童装快闪", "公众 3,000 人", "童装 + 国潮"],
        ["10", "2027-02", "国际设计师论坛", "现场 200 人", "海外设计师"],
        ["11", "2027-03", "春季产业开放周", "系列 5 天", "公众/媒体/高校"],
        ["12", "2027-03", "童装家庭嘉年华（春）", "公众 5,000 人", "童装家庭"],
    ]
    add_table(s, Inches(0.5), Inches(1.2), Inches(8.0), Inches(5.5),
              header, rows,
              col_widths=[Inches(0.4), Inches(1.0), Inches(2.6), Inches(1.7), Inches(2.3)])
    add_text(s, Inches(8.8), Inches(1.2), Inches(4.5), Inches(0.4),
             "媒体四层矩阵", size=14, bold=True, color=NAVY)
    layers2 = [
        ("权威 · 政策", "央视 / 人民日报 / 新华社 / 解放日报", NAVY),
        ("时尚 · 国际", "VOGUE / ELLE / WWD / BoF / Vogue Business", BLUE),
        ("行业 · 垂类", "服装行业资讯 / 童装童趣 / 母婴行业观察", GOLD),
        ("公众 · 社交", "抖音 / 小红书 / B站 / KOL 100 个", GREEN),
    ]
    for i, (t, d, c) in enumerate(layers2):
        y = Inches(1.65) + Inches(1.20) * i
        add_rect(s, Inches(8.8), y, Inches(4.5), Inches(1.05), fill=CLOUD, line=LINE)
        add_rect(s, Inches(8.8), y, Inches(0.18), Inches(1.05), fill=c)
        add_text(s, Inches(9.05), y + Inches(0.08), Inches(4.2), Inches(0.4),
                 t, size=13, bold=True, color=NAVY)
        add_text(s, Inches(9.05), y + Inches(0.45), Inches(4.2), Inches(0.6),
                 d, size=10, color=INK)

    # ============ 25. Phase 4 扉页 ============
    section_cover(prs, "PHASE 4", "商业条款",
                  "Commercial · 4 Tiers · Service Pack · 3-Year Model",
                  ["任务 7 · 四档收费 + 服务包 + 报价单 + 三年财务"])

    # ============ 26. 四档收费 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 4 · 任务 7",
               page_title="四档收费方案矩阵",
               subtitle="链主月费+政策返还 / 设计师月费 / 电商月费+佣金 / 冠名+对赌")
    cards4 = [
        ("甲档", "链主总部 (1#)", "月费 + 政策返还", NAVY,
         ["起始租金 5.5–7.0 元/㎡·天", "整层 1,500–2,500 ㎡",
          "免租 12–18 个月 · 装补 800–1,200/㎡",
          "区级税收 80% 三年返",
          "楼层冠名 + T 台 + 童装亲子协同"]),
        ("乙档", "设计师工作室 (2#)", "月费 + 基础服务", PURPLE,
         ["起始租金 7.0–8.5 元/㎡·天", "80–500 ㎡ 多层组合",
          "免租 3–6 个月 · 装补 300–500/㎡",
          "面料图书馆 / 3D 打样",
          "设计师集合店 + T 台 5 场/年"]),
        ("丙档", "电商基地 (3#)", "月费 + 服务佣金", BLUE,
         ["直播间 30,000 元/月（30 ㎡）", "大型 80,000 元/月（100 ㎡）",
          "跨境运营房 8.0–9.0 元/㎡·天",
          "服务 GMV 5%–10% 抽成",
          "一键开播 + 跨境合规中心"]),
        ("丁档", "冠名 + 对赌", "低租 + 收益分成", GREEN,
         ["折扣租金 60%–70%", "整楼 / T 台 / 集合店冠名",
          "营收/利润/估值/GMV 对赌",
          "超额营收 5%–10% 返现",
          "园区基金 1%–3% 跟投权"]),
    ]
    cw = Inches(2.95); gx = Inches(0.15); cy = Inches(1.20); ch = Inches(5.4)
    for i, (lvl, name, sub, c, bullets) in enumerate(cards4):
        x = Inches(0.5) + (cw + gx) * i
        add_rect(s, x, cy, cw, ch, fill=CLOUD, line=LINE)
        add_rect(s, x, cy, cw, Inches(0.95), fill=c)
        add_text(s, x, cy + Inches(0.05), cw, Inches(0.4), lvl,
                 size=18, bold=True, color=GOLD if c == NAVY else WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, cy + Inches(0.40), cw, Inches(0.4), name,
                 size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x, cy + Inches(0.65), cw, Inches(0.3), sub,
                 size=10, color=WHITE, align=PP_ALIGN.CENTER)
        for j, b in enumerate(bullets):
            by = cy + Inches(1.10) + Inches(0.65) * j
            add_round(s, x + Inches(0.15), by + Inches(0.18), Inches(0.18), Inches(0.18), "●",
                      fill=c, color=WHITE, size=8, bold=True)
            add_text(s, x + Inches(0.40), by, cw - Inches(0.5), Inches(0.55), b,
                     size=10, color=INK)

    # ============ 27. 服务包矩阵 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 4 · 任务 7",
               page_title="服务包矩阵 · 把租金以外的价值显性化",
               subtitle="6 大类 · 12 项服务 · 按四档差异化配置")
    sv_header = ["服务", "甲 链主", "乙 设计师", "丙 电商", "丁 冠名+对赌"]
    sv_rows = [
        ["楼层 / 整楼冠名", "默认 5 年", "—", "—", "默认（叠加 T 台）"],
        ["T 台主秀场", "终身免费", "5 场次/年优惠", "抽签免费", "冠名 + 终身"],
        ["拍摄棚", "≥ 50 场次/年免费", "≥ 10 场次/年", "8 折", "≥ 100 场次/年免费"],
        ["面料图书馆", "配额 100 万", "配额 50 万", "—", "配额 80 万"],
        ["3D 打样 / 数码印花", "30 单/年免费", "20 单/年免费", "—", "50 单/年免费"],
        ["童装亲子区协同", "数据接口", "—", "应用试点", "数据接口"],
        ["跨境合规中心", "优先", "标准", "默认", "优先"],
        ["数据中台", "全栈 API", "标准", "全栈 API", "全栈 API"],
        ["招聘联运 / 落户", "30 个/年", "10 个/年", "集体面试", "30 个/年"],
        ["公关与媒体", "主秀冠名 + 首发", "设计师奖", "Demo 抽签", "主秀冠名 + 首发"],
        ["集合店货架", "主入口", "集合店货架", "—", "主入口"],
        ["律所 / 财税 / IP", "顶尖年框", "优惠折扣", "标准目录", "顶尖年框"],
    ]
    add_table(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5),
              sv_header, sv_rows,
              col_widths=[Inches(2.6), Inches(2.4), Inches(2.4), Inches(2.4), Inches(2.5)],
              header_size=11, body_size=9, body_align=PP_ALIGN.LEFT)
    add_text(s, Inches(0.5), Inches(6.85), Inches(12), Inches(0.3),
             "把租金以外的服务量化为『积分券』——为客户算清账，为我们沉淀粘性",
             size=10, italic=True, color=GREY)

    # ============ 28. 链主报价单样张 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 4 · 任务 7",
               page_title="链主总部报价单样张（甲档 · 保密）",
               subtitle="致 [链主]：1# 9F 整层 ~2,400 ㎡ · 6 年 · 6.0 元/㎡·天 · 一企一策")
    add_rect(s, Inches(0.5), Inches(1.20), Inches(12.3), Inches(0.55), fill=NAVY)
    add_text(s, Inches(0.7), Inches(1.20), Inches(8), Inches(0.55),
             "SENMA · MODE Hub  |  链主总部报价单  |  CONFIDENTIAL",
             size=14, bold=True, color=GOLD, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(8.5), Inches(1.20), Inches(4.0), Inches(0.55),
             "有效期：30 日", size=11, color=WHITE,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.5), Inches(1.95), Inches(6), Inches(0.4),
             "一、商业条款", size=13, bold=True, color=NAVY)
    biz = [
        ("标的", "1# 9F 整层 ~2,400 ㎡（净高 4.5 m · 玻璃幕墙）"),
        ("租赁期", "6 年（含免租期）"),
        ("起始租金", "6.0 元/㎡·天 · 物业 28 元/㎡·月"),
        ("免租期", "14 个月"),
        ("调价机制", "3 年一调 · 调幅 = max(CPI, 5%)"),
        ("装补", "1,000 元/㎡ · 封顶 240 万元 · 里程碑分期"),
        ("履约保证", "6 个月租金"),
        ("续约", "到期前 12 个月可优先续约 5 年"),
    ]
    for i, (k, v) in enumerate(biz):
        y = Inches(2.40) + Inches(0.36) * i
        add_round(s, Inches(0.55), y + Inches(0.06), Inches(1.3), Inches(0.24),
                  k, fill=BLUE, color=WHITE, size=10, bold=True)
        add_text(s, Inches(1.95), y, Inches(4.7), Inches(0.36), v,
                 size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(6.85), Inches(1.95), Inches(6), Inches(0.4),
             "二、一企一策政策包", size=13, bold=True, color=NAVY)
    pol2 = [
        ("税收返还", "区级留成 80% 三年返 / 50% 后两年返"),
        ("一次性奖励", "[一事一议]"),
        ("人才公寓", "100 套"),
        ("落户绿通", "30 个 / 年"),
        ("T 台主秀场", "终身免费会员"),
        ("拍摄棚", "≥ 50 场次/年免费"),
        ("童装亲子协同", "1F 数据接口"),
        ("品牌权益", "9F 楼层冠名 + 1F 集合店主入口"),
    ]
    for i, (k, v) in enumerate(pol2):
        y = Inches(2.40) + Inches(0.36) * i
        add_round(s, Inches(6.90), y + Inches(0.06), Inches(1.3), Inches(0.24),
                  k, fill=GOLD, color=NAVY, size=10, bold=True)
        add_text(s, Inches(8.30), y, Inches(4.4), Inches(0.36), v,
                 size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    add_rect(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.5), fill=BLUE)
    add_text(s, Inches(0.7), Inches(5.4), Inches(12), Inches(0.5),
             "三、签约里程碑：T+30d 意向条款书 → T+90d 正式合同 → T+150d 进场装修 → T+330d 入驻 + 区领导剪彩",
             size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.95), fill=NAVY)
    add_text(s, Inches(0.7), Inches(6.05), Inches(12), Inches(0.4),
             "森马集团 · SENMA · MODE Hub 项目组",
             size=12, bold=True, color=GOLD)
    add_text(s, Inches(0.7), Inches(6.45), Inches(12), Inches(0.45),
             "联系人：项目总监 [姓名]  |  电话：[xxx]  |  邮箱：[xxx@xxx.com]",
             size=11, color=WHITE)

    # ============ 29. 三年财务 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 4 · 任务 7",
               page_title="三年财务测算 · 关键 KPI",
               subtitle="基于 5 万方综合体 · 出租净 35,000 ㎡ · Y3 EBITDA 转正约 1.35 亿元")
    header = ["科目（万元）", "Y1 2026", "Y2 2027", "Y3 2028"]
    rows = [
        ["租金收入", "2,684", "6,260", "9,168"],
        ["物业费收入", "412", "823", "1,082"],
        ["T 台 + 拍摄 + 直播间", "1,400", "3,800", "6,500"],
        ["服务平台佣金", "200", "800", "2,000"],
        ["童装亲子 / 集合店", "200", "600", "1,200"],
        ["政策返还（净计入）", "200", "1,200", "3,500"],
        ["基金管理费", "0", "800", "1,500"],
        ["总收入", "5,096", "14,283", "24,950"],
        ["总成本", "8,000", "9,900", "11,500"],
        ["EBITDA", "−2,904", "+4,383", "+13,450"],
        ["EBIT", "−5,404", "+1,883", "+10,950"],
    ]
    add_table(s, Inches(0.5), Inches(1.2), Inches(7.0), Inches(4.8),
              header, rows,
              col_widths=[Inches(2.4), Inches(1.5), Inches(1.5), Inches(1.6)])

    kpis = [
        ("入驻率", ["35%", "70%", "92%"], BLUE),
        ("链主签约（累计）", ["1", "2", "3"], NAVY),
        ("设计师工作室", ["15", "35", "55"], GOLD),
        ("EBITDA 利润率", ["−57%", "+31%", "+54%"], GREEN),
    ]
    add_text(s, Inches(7.8), Inches(1.2), Inches(5.5), Inches(0.4),
             "三年关键 KPI", size=14, bold=True, color=NAVY)
    for i, (lbl, vals, c) in enumerate(kpis):
        y = Inches(1.65) + Inches(0.95) * i
        add_text(s, Inches(7.8), y + Inches(0.05), Inches(2.3), Inches(0.4),
                 lbl, size=12, bold=True, color=NAVY)
        for j, v in enumerate(vals):
            x = Inches(10.1) + Inches(1.0) * j
            add_round(s, x, y + Inches(0.02), Inches(0.85), Inches(0.55),
                      v, fill=c, color=WHITE, size=12, bold=True)
            add_text(s, x, y + Inches(0.6), Inches(0.85), Inches(0.25),
                     f"Y{j+1}", size=8, color=GREY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.10), Inches(13), Inches(0.4),
             "敏感性（Y3 EBITDA 万元变化）：入驻 ±10% → −2,800/+1,200 · 租金 ±10% → ±1,800 · 政策兑现 −30% → −1,050 · T台/拍摄/直播 −30% → −2,000",
             size=11, color=INK)

    # ============ 30. Phase 5 扉页 ============
    section_cover(prs, "PHASE 5", "落地推进",
                  "Rollout · 12 Months · 5 → 22 Team",
                  ["任务 8 · 12 个月里程碑 + 团队 + RACI + 风险矩阵"])

    # ============ 31. 12 个月甘特 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 5 · 任务 8",
               page_title="12 个月里程碑甘特图",
               subtitle="M1 团队就位 → M5 9 月时装秀 → M12 入驻率 70%")
    months = [f"M{i+1}" for i in range(12)]
    chart_x = Inches(2.8); chart_y = Inches(1.2); chart_w = Inches(10.0)
    col_w = chart_w / 12.0
    for i, m in enumerate(months):
        x = chart_x + col_w * i
        add_rect(s, x, chart_y, col_w, Inches(0.35), fill=NAVY)
        add_text(s, x, chart_y, col_w, Inches(0.35), m,
                 size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tasks = [
        ("立项 + 法人主体",     0, 1, BLUE),
        ("团队组建 (5→22)",     0, 12, NAVY),
        ("样板间 / 视觉",       1, 3, GOLD),
        ("T 台秀场 + 拍摄棚",   1, 5, GREEN),
        ("直播间一期",          1, 5, PURPLE),
        ("品牌官网 + CRM",      1, 3, BLUE),
        ("300 家库 + 中介",     2, 3, NAVY),
        ("首轮政府汇报",        1, 3, GOLD),
        ("政策包 v0 → v1",      3, 5, BLUE),
        ("链主 TOP5 接触",      1, 9, NAVY),
        ("链主首份 Term Sheet", 2, 2, GOLD),
        ("设计师 30 家签约",    4, 4, BLUE),
        ("5 月品牌发布会",      4, 1, GREEN),
        ("6 月童装嘉年华",      5, 1, GOLD),
        ("9 月时装秀（旗舰）",  8, 1, RED),
        ("MODE Hub 10 条",      8, 2, GOLD),
        ("入驻率 50%",          10, 1, GREEN),
        ("入驻率 70%",          11, 1, GREEN),
    ]
    row_h = Inches(0.26)
    for i, (lbl, start, dur, c) in enumerate(tasks):
        y = chart_y + Inches(0.4) + (row_h + Inches(0.03)) * i
        add_text(s, Inches(0.4), y, Inches(2.35), row_h, lbl,
                 size=9, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, chart_x, y, chart_w, row_h, fill=CLOUD)
        bx = chart_x + col_w * start
        bw = col_w * dur
        add_rect(s, bx + Emu(20000), y + Emu(20000), bw - Emu(40000), row_h - Emu(40000), fill=c)

    # ============ 32. 团队 + RACI ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 5 · 任务 8",
               page_title="5 人核心团队 → 22 人稳态 + RACI 分工",
               subtitle="启动 M1–3：5 人 / 扩张 M4–6：14 人 / 稳态 M7–12：22 人")
    add_text(s, Inches(0.5), Inches(1.15), Inches(6), Inches(0.4),
             "5 人核心团队（启动期）", size=14, bold=True, color=NAVY)
    roles = [
        ("项目总监 GM", "全面统筹 / 链主谈判 / 重大客户", NAVY),
        ("BD 总监", "政府关系 / 行业协会 / 品牌资源", BLUE),
        ("招商总监", "渠道 / 漏斗 / 生态签约", GOLD),
        ("运营总监", "物业 / IT / T 台 / 拍摄棚 / 直播间", GREEN),
        ("时尚活动总监", "时装秀 / 媒体 / 设计师 / 童装嘉年华", PURPLE),
    ]
    for i, (t, d, c) in enumerate(roles):
        y = Inches(1.55) + Inches(0.85) * i
        add_rect(s, Inches(0.5), y, Inches(6.0), Inches(0.75), fill=CLOUD, line=LINE)
        add_rect(s, Inches(0.5), y, Inches(0.18), Inches(0.75), fill=c)
        add_text(s, Inches(0.7), y + Inches(0.04), Inches(5.7), Inches(0.4),
                 t, size=14, bold=True, color=NAVY)
        add_text(s, Inches(0.7), y + Inches(0.40), Inches(5.7), Inches(0.35),
                 d, size=11, color=INK)

    add_text(s, Inches(0.5), Inches(6.05), Inches(6), Inches(0.4),
             "扩编节奏：5 → 14 → 22 人", size=12, bold=True, color=NAVY)
    add_text(s, Inches(0.5), Inches(6.40), Inches(6), Inches(0.6),
             "招商 1→5→8 · BD 1→2→2 · 运营 1→3→5 · 市场/活动 1→2→4 · 法务 0→1→1 · HR/财务 0→1→2",
             size=10, color=INK)

    add_text(s, Inches(7.0), Inches(1.15), Inches(6), Inches(0.4),
             "RACI 关键事项分工", size=14, bold=True, color=NAVY)
    header = ["关键事项", "总监", "BD", "招商", "运营", "活动", "董事长"]
    rows = [
        ["链主谈判", "R", "C", "C", "I", "I", "A"],
        ["政府汇报 / 政策包", "C", "R", "I", "I", "I", "A"],
        ["中介 / 漏斗例会", "A", "I", "R", "I", "I", "I"],
        ["9 月时装秀 / T 台", "A", "C", "C", "C", "R", "C"],
        ["合同 / 对赌", "A", "C", "C", "I", "I", "A"],
        ["童装亲子区 / 集合店", "A", "C", "C", "C", "R", "A"],
        ["直播 / 跨境合规", "C", "C", "C", "C", "I", "A"],
    ]
    add_table(s, Inches(7.0), Inches(1.55), Inches(6.0), Inches(4.0),
              header, rows,
              col_widths=[Inches(2.0), Inches(0.6), Inches(0.6), Inches(0.7),
                          Inches(0.7), Inches(0.7), Inches(0.7)],
              header_size=10, body_size=10, body_align=PP_ALIGN.CENTER)
    add_text(s, Inches(7.0), Inches(5.65), Inches(6.0), Inches(0.5),
             "R = Responsible · A = Accountable · C = Consulted · I = Informed",
             size=10, italic=True, color=GREY)

    # ============ 33. 风险矩阵 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="Phase 5 · 任务 8",
               page_title="风险矩阵 · 高频高影响项的对冲方案",
               subtitle="按概率 × 影响排序，每项配对应『对冲手段』和『启动条件』")
    rk_header = ["#", "风险", "概率", "影响", "对冲手段", "启动条件"]
    rk_rows = [
        ["R1", "链主推进慢于预期", "中", "高", "启动备选名单（UR / MO&Co. / Edition / SHIATZY 等）+ 链主分阶段意向", "TOP5 中 ≤1 进入 Term Sheet @ M6"],
        ["R2", "T 台 / 拍摄棚 报批延迟", "中", "高", "第三方代办 + 与文化执法 / 公安预先沟通 + 中汽研合作", "M4 仍未取得审批"],
        ["R3", "政策返还兑现延迟", "中", "中", "改约定 + 财政确认函分期 + 集团担保过渡", "区财政书面确认延误 > 60 天"],
        ["R4", "直播 GMV 下滑 / 监管收紧", "中", "高", "转向跨境 + 海外仓；引入 SHEIN/TEMU 流量", "全行业 GMV 同比 -20%"],
        ["R5", "团队建设速度", "中", "中", "启动期外包（PR/HR/法务）+ 头部猎头授权", "招聘缺口 > 30% @ M3"],
        ["R6", "资金压力", "中", "中", "Y2 起 REITs 化探索 + 集团信用借款", "现金安全垫 < 6 个月"],
        ["R7", "设计师品牌存活率低 (<30%)", "中", "中", "设计师孵化基金 + 集合店扶持 + 跨境出海", "12 个月续约率 < 50%"],
        ["R8", "童模 / 直播内容违规", "低", "高", "严守国家规定 + 法务声明库 + 内审 SOP", "行业重大事件 / 监管整顿"],
        ["R9", "施工 / 装修延期", "中", "中", "工程总包 + 关键节点违约金", "里程碑滞后 > 30 天"],
        ["R10", "跨境政策不确定", "中", "中", "多渠道（TikTok/Shopee/Shopify）+ 多区域海外仓", "跨境政策重大调整"],
    ]
    add_table(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5),
              rk_header, rk_rows,
              col_widths=[Inches(0.5), Inches(2.4), Inches(0.7), Inches(0.7),
                          Inches(4.5), Inches(3.5)],
              header_size=11, body_size=9)
    add_text(s, Inches(0.5), Inches(6.85), Inches(13), Inches(0.3),
             "风险委员会按月评审；任何风险升至『红』状态，72 小时内升级至集团董事长 + 项目总监 + 法务三方专项会",
             size=10, italic=True, color=GREY)

    # ============ 34. 投决建议 ============
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="07 · 投决建议",
               page_title="投决建议 · 里程碑承诺",
               subtitle="资源到位的前提下，T+30 / 90 / 180 三个节点的硬承诺")
    nodes2 = [
        ("T + 30 天", "团队就位 + 政府首轮汇报",
         ["5 人核心入职", "区投促办 + 副区长 + 市经信委首轮汇报", "5 家中介签约"], BLUE),
        ("T + 90 天", "链主首份 Term Sheet",
         ["TOP5 全部进入会面阶段", "≥ 1 家链主签 Term Sheet", "300 家库 + CRM 上线"], GOLD),
        ("T + 180 天", "9 月时装秀 + 政策包定稿",
         ["时装秀 3,000 人次 + 5 家签约", "MODE Hub 10 条发布", "签约 ≥ 40 家 / 入驻率 ≥ 30%"], GREEN),
    ]
    for i, (t, sub, bullets, c) in enumerate(nodes2):
        x = Inches(0.5) + Inches(4.30) * i
        add_rect(s, x, Inches(1.2), Inches(4.10), Inches(4.5), fill=CLOUD, line=LINE)
        add_rect(s, x, Inches(1.2), Inches(4.10), Inches(0.95), fill=c)
        add_text(s, x, Inches(1.25), Inches(4.10), Inches(0.5), t,
                 size=22, bold=True, color=NAVY if c == GOLD else WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(1.70), Inches(4.10), Inches(0.4), sub,
                 size=12, color=NAVY if c == GOLD else WHITE, align=PP_ALIGN.CENTER)
        for j, b in enumerate(bullets):
            y = Inches(2.35) + Inches(0.65) * j
            add_round(s, x + Inches(0.20), y + Inches(0.15), Inches(0.25), Inches(0.25),
                      str(j + 1), fill=c, color=WHITE if c != GOLD else NAVY,
                      size=10, bold=True)
            add_text(s, x + Inches(0.55), y + Inches(0.05), Inches(3.40), Inches(0.55),
                     b, size=12, color=INK)
    add_rect(s, Inches(0.5), Inches(5.95), Inches(12.3), Inches(1.0), fill=NAVY)
    add_text(s, Inches(0.7), Inches(6.0), Inches(12.0), Inches(0.45),
             "请董事会审议", size=14, bold=True, color=GOLD)
    add_text(s, Inches(0.7), Inches(6.40), Inches(12.0), Inches(0.55),
             "① 启动预算（Y0–Y1）  ② 5 人核心团队招聘授权  ③ 区一企一策政府专班 ④ 9 月时装秀预算  ⑤ 链主谈判授权区间",
             size=12, color=WHITE)

    # ============ 35. Q&A ============
    s = new_slide(prs)
    add_rect(s, 0, 0, sw, sh, fill=NAVY)
    add_rect(s, 0, Inches(3.6), sw, Emu(40000), fill=GOLD)
    add_text(s, Inches(0.8), Inches(2.3), Inches(11), Inches(1.4),
             "SENMA · MODE Hub", size=44, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), Inches(3.8), Inches(11), Inches(0.7),
             "让中国时尚走向世界",
             size=22, color=WHITE, italic=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), Inches(5.0), Inches(11), Inches(0.5),
             "Q & A · 谢 谢", size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), Inches(6.4), Inches(11), Inches(0.4),
             "森马集团 · SENMA · MODE Hub 项目组",
             size=12, color=CLOUD, align=PP_ALIGN.CENTER)

    # 回填总页数
    total = len(SLIDES)
    for sl in SLIDES:
        for shape in sl.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text and run.text.strip().endswith(" / 0"):
                        run.text = run.text.replace(" / 0", f" / {total}")

    out = Path(__file__).resolve().parent.parent / "docs" / "deck" / \
        "SENMA-MODE-Hub-招商方案.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"✓ Deck written: {out}  ({total} slides)")
    return out


if __name__ == "__main__":
    build()
