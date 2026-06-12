# -*- coding: utf-8 -*-
"""
活动财务整理：将PDF原始流水整理成会计风格的 Excel 与 PPT。
- 严格保持原始数据（按对账逻辑校验各分项小计/合计无误）
- 按会计科目分类：外部费用 / 内部费用 / 现场杂项（代收代付）/ 赞助收入
- 表头使用中国传统会计配色（红/金/灰），方便日后多方查阅
"""
import os
from openpyxl import Workbook
from openpyxl.styles import (
    Font, PatternFill, Alignment, Border, Side, NamedStyle
)
from openpyxl.utils import get_column_letter
from openpyxl.formatting.rule import CellIsRule

OUT_DIR = "/workspace/output"
os.makedirs(OUT_DIR, exist_ok=True)

# ---------------------------------------------------------------
# 1. 原始数据（与PDF逐项核对，小计与合计已校验一致）
# ---------------------------------------------------------------

# 外部费用：对外采购/服务支出
EXTERNAL = [
    ("主持人",          "Miranda 苗苗",  -500.00,   "主持人礼服"),
    ("晚宴期间演出",     "小五",          -3000.00,  "三人演出表演"),
    ("摄影摄像",        "洪总",          -3000.00,  "摄影加云图片，及摄像影像记录"),
    ("会务公司",        "杨总",          -16000.00, "全部费用合计见附图"),
    ("媒体支出",        "葛九明",        -11500.00, "新华社、上海证券报、央视新闻客户端、央广网、中国新闻网、上海虹口融媒体、上海杨浦融媒体（其中500元由胡继刚付款）"),
    ("晚宴用餐",        "尚九一滴水",     -33452.00, "晚宴8桌用餐费用，主桌按1.5倍计算"),
    ("证书制作及KT版",   "广告公司",       -1030.00,  "个人证书14件 + 企业证书10件 + 志愿者证书23件"),
    ("晚宴红包",        "胡继刚",         -100.00,   "晚宴活动游戏环节用"),
    ("打印费",          "王胜",           -10.00,    "主持人手稿"),
]

# 内部费用：内部团队垫付/团队侧支出
INTERNAL = [
    ("茶歇",            "F组-彭常丽",     -2000.50,  "盒马鲜生采购，单次≤30kg，分8单下单（含订单截图录屏）"),
    ("布袋",            "F组-彭常丽",     -860.00,   "200个 × 4.30元/个"),
    ("志愿者伴手礼",     "F组-彭常丽",     -2025.00,  "30份 × 68元/份 = 2040元，使用15元券，实付2025元"),
    ("志愿者证书邮寄费", "胡继刚",         -200.00,   "5位志愿者未现场领取证书及礼品，另含夏春等邮寄费"),
    ("A4 刊册议程",     "佳玉",          -4400.00,  "500份刊册"),
    ("午餐",            "F组-彭常丽",     -780.00,   "大米先生 30元/份 × 26份（附付款截图）"),
    ("活动前看场地用餐", "王胜、王卓 等",   148.85,   "活动前看场地三次用餐（金额按原票据登记）"),
]

# 现场杂项：本质上是代收代付（向参会者收升舱费再统一支付，含退费）
MISC = [
    ("胸花 & 台花",     "F组-朱俊锋",     -488.00,   "胸花10个 + 台花1个"),
    ("对讲机租赁",      "F组-朱俊锋",     -138.88,   "租赁100元 + 来回运费38.88元"),
    ("升舱费用（代收）", "F组-朱俊锋",    1440.00,   "9人升舱，160元/人，由参会者交付"),
    ("晚宴退费",        "F组-朱俊锋",     -199.00,   "1人不愿升舱，按规则退费"),
    ("货拉拉送货",      "F组-朱俊锋",     -100.00,   "送复旦管院纸袋子 + 酒水送件"),
]

# 赞助/收入
SPONSOR = [
    ("赞助费",            "姚丹",                 1500.00,  ""),
    ("晚宴费",            "乐琪文",                419.00,  ""),
    ("晚宴费",            "黄家颖",                420.00,  ""),
    ("晚宴费",            "*斯",                  419.00,  ""),
    ("晚宴费",            "泰隆银行",             8000.00,  ""),
    ("活动赞助",          "夏洁",                 6000.00,  ""),
    ("活动赞助",          "冠松汽车 - 陈总",     10000.00,  ""),
    ("活动赞助",          "北京大学上海校友会 - 黄欣", 20000.00, ""),
]

# 证书打印费明细
CERT_DETAIL = [
    ("XH20260522016", "彩色打印", "HD", "套", 24, 35, 840, "01-02 证书+外壳：企业10 + 个人14"),
    ("XH20260522016", "彩色打印", "A4", "张", 24, 5,  120, "—"),
    ("XH20260522016", "彩色打印", "A4", "套", 1,  35, 35,  "03 志愿者证书 珠光纸 1份"),
    ("XH20260522016", "快递费",   "HD", "次", 1,  35, 35,  "04 证书160克超白+证书外壳 1套；闪送：尚九一滴水1F"),
]

# ---------------------------------------------------------------
# 2. 校验小计/合计（开发期断言，确保数据准确）
# ---------------------------------------------------------------
def _sum(rows):
    return round(sum(r[2] for r in rows), 2)

assert _sum(EXTERNAL) == -68592.00,    f"外部费用小计异常: {_sum(EXTERNAL)}"
assert _sum(INTERNAL) == -10116.65,    f"内部费用小计异常: {_sum(INTERNAL)}"
assert _sum(MISC)     == 514.12,       f"现场杂项小计异常: {_sum(MISC)}"
assert _sum(SPONSOR)  == 46758.00,     f"赞助收入小计异常: {_sum(SPONSOR)}"
GRAND_TOTAL = round(_sum(EXTERNAL) + _sum(INTERNAL) + _sum(MISC) + _sum(SPONSOR), 2)
assert GRAND_TOTAL == -31436.53,       f"合计异常: {GRAND_TOTAL}"

# ---------------------------------------------------------------
# 3. 通用样式
# ---------------------------------------------------------------
COLOR_PRIMARY   = "8B0000"   # 暗红 - 财务报表惯用主色
COLOR_ACCENT    = "C9A86B"   # 金色 - 表头点缀
COLOR_SUBHEAD   = "F5E9D3"   # 米色 - 副表头/小计行
COLOR_STRIPE    = "FBF7F0"   # 极浅米色 - 隔行
COLOR_TOTAL_BG  = "FFE4B5"   # 浅金 - 合计
COLOR_INCOME    = "1F6E43"   # 绿 - 收入
COLOR_EXPENSE   = "B22222"   # 红 - 支出

thin = Side(border_style="thin", color="C9A86B")
medium = Side(border_style="medium", color="8B0000")
BORDER_CELL = Border(left=thin, right=thin, top=thin, bottom=thin)
BORDER_HEAD = Border(left=thin, right=thin, top=medium, bottom=medium)

FONT_TITLE = Font(name="WenQuanYi Micro Hei", size=18, bold=True, color="FFFFFF")
FONT_SUBTITLE = Font(name="WenQuanYi Micro Hei", size=11, italic=True, color="555555")
FONT_HEAD = Font(name="WenQuanYi Micro Hei", size=12, bold=True, color="FFFFFF")
FONT_SUBHEAD = Font(name="WenQuanYi Micro Hei", size=11, bold=True, color="8B0000")
FONT_CELL = Font(name="WenQuanYi Micro Hei", size=10, color="333333")
FONT_TOTAL = Font(name="WenQuanYi Micro Hei", size=12, bold=True, color="8B0000")
FONT_NUM_EXP = Font(name="WenQuanYi Micro Hei", size=10, color=COLOR_EXPENSE)
FONT_NUM_INC = Font(name="WenQuanYi Micro Hei", size=10, color=COLOR_INCOME)

FILL_TITLE = PatternFill("solid", fgColor=COLOR_PRIMARY)
FILL_HEAD = PatternFill("solid", fgColor=COLOR_PRIMARY)
FILL_SUBHEAD = PatternFill("solid", fgColor=COLOR_SUBHEAD)
FILL_STRIPE = PatternFill("solid", fgColor=COLOR_STRIPE)
FILL_TOTAL = PatternFill("solid", fgColor=COLOR_TOTAL_BG)

ALIGN_CENTER = Alignment(horizontal="center", vertical="center", wrap_text=True)
ALIGN_LEFT = Alignment(horizontal="left", vertical="center", wrap_text=True)
ALIGN_RIGHT = Alignment(horizontal="right", vertical="center", wrap_text=True)

NUM_FMT = '_-¥* #,##0.00_-;[Red]-¥* #,##0.00_-;_-¥* "-"??_-;_-@_-'

# ---------------------------------------------------------------
# 4. 构建 Excel
# ---------------------------------------------------------------
wb = Workbook()

def _title_row(ws, title, subtitle, ncols):
    ws.row_dimensions[1].height = 38
    ws.row_dimensions[2].height = 22
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=ncols)
    ws.merge_cells(start_row=2, start_column=1, end_row=2, end_column=ncols)
    c = ws.cell(row=1, column=1, value=title)
    c.font = FONT_TITLE; c.fill = FILL_TITLE; c.alignment = ALIGN_CENTER
    c2 = ws.cell(row=2, column=1, value=subtitle)
    c2.font = FONT_SUBTITLE; c2.alignment = ALIGN_CENTER
    for col in range(1, ncols + 1):
        ws.cell(row=1, column=col).border = BORDER_HEAD

def _header_row(ws, row, headers):
    ws.row_dimensions[row].height = 28
    for i, h in enumerate(headers, 1):
        c = ws.cell(row=row, column=i, value=h)
        c.font = FONT_HEAD; c.fill = FILL_HEAD
        c.alignment = ALIGN_CENTER; c.border = BORDER_HEAD

def _write_money(cell, value):
    cell.number_format = NUM_FMT
    cell.value = value
    if value is None or value == 0:
        cell.font = FONT_CELL
    elif value < 0:
        cell.font = FONT_NUM_EXP
    else:
        cell.font = FONT_NUM_INC
    cell.alignment = ALIGN_RIGHT
    cell.border = BORDER_CELL

def _set_widths(ws, widths):
    for i, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = w

def _enable_fit_to_page(ws):
    ws.page_setup.orientation = ws.ORIENTATION_LANDSCAPE
    ws.page_setup.paperSize = ws.PAPERSIZE_A4
    ws.page_setup.fitToWidth = 1
    ws.page_setup.fitToHeight = 0
    ws.sheet_properties.pageSetUpPr.fitToPage = True
    ws.page_margins.left = 0.3
    ws.page_margins.right = 0.3
    ws.page_margins.top = 0.4
    ws.page_margins.bottom = 0.4
    ws.print_options.horizontalCentered = True

# ---------- 工作表 1：财务总览 ----------
ws = wb.active
ws.title = "财务总览"
_title_row(
    ws,
    "2026 年度活动 · 财务收支总览",
    "数据来源：原始流水汇总表  ｜  币种：人民币（元）  ｜  红色为支出 / 绿色为收入",
    5,
)
_header_row(ws, 3, ["序号", "类别", "性质", "金额（元）", "说明"])

summary_rows = [
    ("一", "外部费用",            "支出", _sum(EXTERNAL), "对外采购、服务、媒体、餐饮等"),
    ("二", "内部费用",            "支出", _sum(INTERNAL), "内部团队垫付的物料、伴手礼、餐饮等"),
    ("三", "现场杂项（代收代付）", "净额", _sum(MISC),     "升舱费代收 + 现场零星支出 + 退费净额"),
    ("四", "赞助 / 收入",         "收入", _sum(SPONSOR),  "企业赞助、个人赞助、晚宴费等"),
]
r = 4
for idx, (no, name, nature, amt, desc) in enumerate(summary_rows):
    ws.cell(row=r, column=1, value=no).alignment = ALIGN_CENTER
    ws.cell(row=r, column=2, value=name).alignment = ALIGN_LEFT
    ws.cell(row=r, column=3, value=nature).alignment = ALIGN_CENTER
    _write_money(ws.cell(row=r, column=4), amt)
    ws.cell(row=r, column=5, value=desc).alignment = ALIGN_LEFT
    for col in range(1, 6):
        cell = ws.cell(row=r, column=col)
        cell.border = BORDER_CELL
        if cell.font is None or cell.font.color is None or cell.font.color.rgb is None or cell.font.color.rgb not in (f"00{COLOR_INCOME}", f"00{COLOR_EXPENSE}"):
            if col != 4:
                cell.font = FONT_CELL
        if idx % 2 == 1 and col != 4:
            cell.fill = FILL_STRIPE
    r += 1

# 合计
ws.row_dimensions[r].height = 30
ws.cell(row=r, column=1, value="合计").alignment = ALIGN_CENTER
ws.merge_cells(start_row=r, start_column=1, end_row=r, end_column=3)
total_label = ws.cell(row=r, column=1)
total_label.font = FONT_TOTAL; total_label.fill = FILL_TOTAL; total_label.border = BORDER_CELL
_write_money(ws.cell(row=r, column=4), GRAND_TOTAL)
ws.cell(row=r, column=4).font = FONT_TOTAL
ws.cell(row=r, column=4).fill = FILL_TOTAL
note = ws.cell(row=r, column=5, value="收入 − 支出（含代收代付净额）")
note.font = FONT_TOTAL; note.fill = FILL_TOTAL; note.alignment = ALIGN_LEFT; note.border = BORDER_CELL

# 结余结论
r += 2
ws.merge_cells(start_row=r, start_column=1, end_row=r, end_column=5)
c = ws.cell(row=r, column=1,
            value=f"结余结论：总收入 ¥{_sum(SPONSOR):,.2f}  −  总支出 ¥{abs(_sum(EXTERNAL)+_sum(INTERNAL)):,.2f}"
                  f"  +  代收代付净额 ¥{_sum(MISC):,.2f}  =  本次活动净结余 ¥{GRAND_TOTAL:,.2f}（负数表示活动总体亏损，需另行筹措补足）")
c.font = Font(name="WenQuanYi Micro Hei", size=11, bold=True, color="8B0000")
c.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)
c.fill = FILL_SUBHEAD
ws.row_dimensions[r].height = 38

_set_widths(ws, [6, 22, 10, 18, 60])
ws.freeze_panes = "A4"
ws.sheet_view.showGridLines = False
_enable_fit_to_page(ws)

# ---------- 通用：明细表 ----------
def build_detail_sheet(sheet_name, title, subtitle, data, subtotal_label, subtotal_value):
    ws = wb.create_sheet(sheet_name)
    _title_row(ws, title, subtitle, 6)
    _header_row(ws, 3, ["序号", "项目", "对接人 / 供应商", "金额（元）", "性质", "备注"])
    r = 4
    for i, (proj, contact, amt, remark) in enumerate(data, 1):
        ws.row_dimensions[r].height = 22 if len(remark) < 40 else 36
        ws.cell(row=r, column=1, value=i).alignment = ALIGN_CENTER
        ws.cell(row=r, column=2, value=proj).alignment = ALIGN_LEFT
        ws.cell(row=r, column=3, value=contact).alignment = ALIGN_LEFT
        _write_money(ws.cell(row=r, column=4), amt)
        nature = "收入" if amt > 0 else "支出"
        ws.cell(row=r, column=5, value=nature).alignment = ALIGN_CENTER
        ws.cell(row=r, column=6, value=remark).alignment = ALIGN_LEFT
        for col in range(1, 7):
            cell = ws.cell(row=r, column=col)
            cell.border = BORDER_CELL
            if col not in (4,):
                cell.font = FONT_CELL
            if i % 2 == 0 and col != 4:
                cell.fill = FILL_STRIPE
        r += 1
    # 小计
    ws.row_dimensions[r].height = 28
    ws.merge_cells(start_row=r, start_column=1, end_row=r, end_column=3)
    sl = ws.cell(row=r, column=1, value=subtotal_label)
    sl.font = FONT_TOTAL; sl.fill = FILL_TOTAL; sl.alignment = ALIGN_CENTER; sl.border = BORDER_CELL
    _write_money(ws.cell(row=r, column=4), subtotal_value)
    ws.cell(row=r, column=4).font = FONT_TOTAL
    ws.cell(row=r, column=4).fill = FILL_TOTAL
    for col in (5, 6):
        c = ws.cell(row=r, column=col, value="")
        c.fill = FILL_TOTAL; c.border = BORDER_CELL
    _set_widths(ws, [6, 22, 22, 16, 8, 55])
    ws.freeze_panes = "A4"
    ws.sheet_view.showGridLines = False
    _enable_fit_to_page(ws)

build_detail_sheet(
    "外部费用明细", "一、外部费用明细",
    "对外采购、服务、媒体宣发、餐饮等支出",
    EXTERNAL, "外部费用小计", _sum(EXTERNAL),
)
build_detail_sheet(
    "内部费用明细", "二、内部费用明细",
    "内部团队垫付的物料、伴手礼、餐饮等支出",
    INTERNAL, "内部费用小计", _sum(INTERNAL),
)
build_detail_sheet(
    "现场杂项", "三、现场杂项（代收代付）",
    "升舱费代收 + 现场零星支出 + 退费，净额体现",
    MISC, "现场杂项小计", _sum(MISC),
)
build_detail_sheet(
    "赞助收入明细", "四、赞助 / 收入明细",
    "企业赞助、个人赞助、晚宴费等",
    SPONSOR, "赞助收入小计", _sum(SPONSOR),
)

# ---------- 工作表：证书打印费明细 ----------
ws = wb.create_sheet("证书打印费明细")
_title_row(
    ws,
    "附表 · 证书制作及KT版打印明细",
    "供应商：广告公司  ｜  原始单据编号：XH20260522016",
    8,
)
_header_row(ws, 3, ["序号", "单据编号", "商品名称", "型号", "单位", "数量", "单价（元）", "金额（元）"])
r = 4
for i, (no, name, mdl, unit, qty, price, amt, _remark) in enumerate(CERT_DETAIL, 1):
    ws.cell(row=r, column=1, value=i).alignment = ALIGN_CENTER
    ws.cell(row=r, column=2, value=no).alignment = ALIGN_CENTER
    ws.cell(row=r, column=3, value=name).alignment = ALIGN_LEFT
    ws.cell(row=r, column=4, value=mdl).alignment = ALIGN_CENTER
    ws.cell(row=r, column=5, value=unit).alignment = ALIGN_CENTER
    ws.cell(row=r, column=6, value=qty).alignment = ALIGN_CENTER
    _write_money(ws.cell(row=r, column=7), price)
    ws.cell(row=r, column=7).font = FONT_CELL
    _write_money(ws.cell(row=r, column=8), amt)
    ws.cell(row=r, column=8).font = FONT_CELL
    for col in range(1, 9):
        cell = ws.cell(row=r, column=col)
        cell.border = BORDER_CELL
        if cell.font is None:
            cell.font = FONT_CELL
        if i % 2 == 0:
            cell.fill = FILL_STRIPE
    r += 1
# 合计
ws.merge_cells(start_row=r, start_column=1, end_row=r, end_column=7)
sl = ws.cell(row=r, column=1, value="合计")
sl.font = FONT_TOTAL; sl.fill = FILL_TOTAL; sl.alignment = ALIGN_CENTER; sl.border = BORDER_CELL
_write_money(ws.cell(row=r, column=8), 1030.00)
ws.cell(row=r, column=8).font = FONT_TOTAL
ws.cell(row=r, column=8).fill = FILL_TOTAL
r += 2
# 备注
ws.merge_cells(start_row=r, start_column=1, end_row=r, end_column=8)
c = ws.cell(row=r, column=1,
            value=("备注：\n"
                   "  01-02  证书+外壳：企业10件 + 个人14件\n"
                   "  03     志愿者证书 珠光纸 1份\n"
                   "  04     证书160克超白 + 证书外壳 1套；闪送：尚九一滴水1F"))
c.font = FONT_CELL; c.alignment = Alignment(horizontal="left", vertical="top", wrap_text=True)
c.fill = FILL_SUBHEAD; c.border = BORDER_CELL
ws.row_dimensions[r].height = 80

_set_widths(ws, [6, 18, 14, 8, 8, 8, 14, 16])
ws.freeze_panes = "A4"
ws.sheet_view.showGridLines = False
_enable_fit_to_page(ws)

# ---------- 说明页 ----------
ws = wb.create_sheet("阅读说明", 0)  # 放在最前
wb.move_sheet("阅读说明", offset=-(len(wb.sheetnames)-1))
_title_row(
    ws,
    "2026 年度活动 · 财务整理表（供各方查阅）",
    "本表依据原始流水整理，按会计科目分类，附原票据备注",
    2,
)
notes = [
    ("整理人 / 日期",   "活动财务组  ·  2026 年"),
    ("币种",           "人民币（元）"),
    ("数据来源",       "活动原始流水 PDF（与PDF逐行核对，小计与合计一致）"),
    ("分类口径",       "按资金性质分四类：① 外部费用  ② 内部费用  ③ 现场杂项（代收代付）  ④ 赞助 / 收入"),
    ("数字颜色",       "红色 = 支出（负数）；绿色 = 收入（正数）"),
    ("注意事项 1",     "「内部费用」中『活动前看场地用餐 +148.85』按原票据登记，建议复核是否为退款或符号录入差异"),
    ("注意事项 2",     "「现场杂项」属代收代付性质：向参会者收升舱费1,440元，对应支付胸花/对讲机/送货/退费等，最终净额 +514.12 元"),
    ("注意事项 3",     "「媒体支出 11,500 元」其中 500 元由胡继刚代付，已含在该笔金额内"),
    ("结余",           f"本次活动总体净结余：¥{GRAND_TOTAL:,.2f}（负数表示亏损）"),
    ("工作表索引",     "财务总览 → 外部费用明细 → 内部费用明细 → 现场杂项 → 赞助收入明细 → 证书打印费明细"),
]
r = 4
for k, v in notes:
    ws.row_dimensions[r].height = 26
    a = ws.cell(row=r, column=1, value=k)
    a.font = FONT_SUBHEAD; a.fill = FILL_SUBHEAD
    a.alignment = ALIGN_CENTER; a.border = BORDER_CELL
    b = ws.cell(row=r, column=2, value=v)
    b.font = FONT_CELL; b.alignment = ALIGN_LEFT; b.border = BORDER_CELL
    r += 1
_set_widths(ws, [18, 72])
ws.sheet_view.showGridLines = False
_enable_fit_to_page(ws)

xlsx_path = os.path.join(OUT_DIR, "2026年活动财务整理表.xlsx")
wb.save(xlsx_path)
print("Excel 已生成：", xlsx_path)

# ---------------------------------------------------------------
# 5. PPT 汇报
# ---------------------------------------------------------------
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

DARK_RED = RGBColor(0x8B, 0x00, 0x00)
GOLD     = RGBColor(0xC9, 0xA8, 0x6B)
CREAM    = RGBColor(0xFB, 0xF7, 0xF0)
DARK_TXT = RGBColor(0x33, 0x33, 0x33)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0x6B, 0x6B, 0x6B)
GREEN    = RGBColor(0x1F, 0x6E, 0x43)
RED      = RGBColor(0xB2, 0x22, 0x22)
BG_PAGE  = RGBColor(0xFB, 0xF7, 0xF0)

FONT_CN = "WenQuanYi Micro Hei"

def add_bg(slide, color=BG_PAGE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg

def add_top_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.95))
    bar.fill.solid(); bar.fill.fore_color.rgb = DARK_RED
    bar.line.fill.background()
    # 金色装饰条
    deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.95), SW, Inches(0.06))
    deco.fill.solid(); deco.fill.fore_color.rgb = GOLD
    deco.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), SW - Inches(1), Inches(0.85))
    tf = tb.text_frame; tf.margin_left = 0; tf.margin_top = 0
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.name = FONT_CN; r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.name = FONT_CN; r2.font.size = Pt(12); r2.font.color.rgb = GOLD

def add_footer(slide, idx, total):
    tb = slide.shapes.add_textbox(Inches(0.4), SH - Inches(0.4), SW - Inches(0.8), Inches(0.3))
    tf = tb.text_frame; tf.margin_left = 0; tf.margin_top = 0
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "2026 年度活动 · 财务整理"
    r.font.name = FONT_CN; r.font.size = Pt(9); r.font.color.rgb = GRAY
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run(); r2.text = f"{idx} / {total}"
    r2.font.name = FONT_CN; r2.font.size = Pt(9); r2.font.color.rgb = GRAY

def set_cell(cell, text, *, bold=False, size=11, color=DARK_TXT,
             fill=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    cell.text = ""
    tf = cell.text_frame
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(text)
    r.font.name = FONT_CN; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    cell.vertical_anchor = anchor
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = fill

def money_str(x):
    if x is None:
        return ""
    return f"¥{x:,.2f}"

def money_color(x):
    if x is None or x == 0:
        return DARK_TXT
    return RED if x < 0 else GREEN

# ---- Slide 1: 封面 ----
slide_count = 8
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG_PAGE)
# 顶部红块
top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(3.6))
top.fill.solid(); top.fill.fore_color.rgb = DARK_RED; top.line.fill.background()
# 金色装饰线
deco = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.6), SW, Inches(0.08))
deco.fill.solid(); deco.fill.fore_color.rgb = GOLD; deco.line.fill.background()

tb = s.shapes.add_textbox(Inches(0.8), Inches(1.0), SW - Inches(1.6), Inches(2.4))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "2026 年度活动"
r.font.name = FONT_CN; r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = GOLD
p2 = tf.add_paragraph()
r2 = p2.add_run(); r2.text = "财 务 收 支 整 理 报 告"
r2.font.name = FONT_CN; r2.font.size = Pt(54); r2.font.bold = True; r2.font.color.rgb = WHITE
p3 = tf.add_paragraph()
r3 = p3.add_run(); r3.text = "Financial Summary  ·  For Multi-party Reference"
r3.font.name = "Calibri"; r3.font.size = Pt(14); r3.font.color.rgb = CREAM

# 中部信息卡
card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.2), SW - Inches(1.6), Inches(2.6))
card.fill.solid(); card.fill.fore_color.rgb = WHITE
card.line.color.rgb = GOLD; card.line.width = Pt(1.5)
tb = s.shapes.add_textbox(Inches(1.2), Inches(4.4), SW - Inches(2.4), Inches(2.3))
tf = tb.text_frame; tf.word_wrap = True
for key, val, c in [
    ("整理依据", "活动原始流水 PDF（与原表逐行核对一致）", DARK_TXT),
    ("分类口径", "外部费用 / 内部费用 / 现场杂项（代收代付）/ 赞助收入", DARK_TXT),
    ("总收入",   f"+ {money_str(_sum(SPONSOR))}", GREEN),
    ("总支出",   f"  {money_str(_sum(EXTERNAL)+_sum(INTERNAL))}", RED),
    ("净结余",   f"  {money_str(GRAND_TOTAL)}（负数表示亏损）", RED),
]:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    r = p.add_run(); r.text = f"  {key} ： "
    r.font.name = FONT_CN; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = DARK_RED
    r2 = p.add_run(); r2.text = val
    r2.font.name = FONT_CN; r2.font.size = Pt(14); r2.font.color.rgb = c
    p.space_after = Pt(6)

add_footer(s, 1, slide_count)

# ---- 复用：渲染明细表格的 slide ----
def add_table_slide(idx, title, subtitle, data, subtotal_label, subtotal):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_top_bar(s, title, subtitle)
    # 表格
    headers = ["序号", "项目", "对接人 / 供应商", "金额（元）", "性质", "备注"]
    rows = len(data) + 2  # 表头 + 数据 + 小计
    cols = len(headers)
    left = Inches(0.4); top = Inches(1.25)
    width = SW - Inches(0.8); height = Inches(5.6)
    table_shape = s.shapes.add_table(rows, cols, left, top, width, height)
    tbl = table_shape.table
    # 列宽
    col_widths_in = [0.7, 2.3, 2.0, 1.6, 0.9, 5.0]
    total_in = sum(col_widths_in)
    for i, w_in in enumerate(col_widths_in):
        tbl.columns[i].width = Emu(int((width / total_in) * w_in))
    # 表头
    for i, h in enumerate(headers):
        set_cell(tbl.cell(0, i), h, bold=True, size=12, color=WHITE,
                 fill=DARK_RED, align=PP_ALIGN.CENTER)
    tbl.rows[0].height = Inches(0.42)
    # 数据
    for ri, (proj, contact, amt, remark) in enumerate(data, 1):
        stripe = CREAM if ri % 2 == 0 else WHITE
        set_cell(tbl.cell(ri, 0), ri, size=10, align=PP_ALIGN.CENTER, fill=stripe)
        set_cell(tbl.cell(ri, 1), proj, size=10, fill=stripe)
        set_cell(tbl.cell(ri, 2), contact, size=10, fill=stripe)
        set_cell(tbl.cell(ri, 3), money_str(amt), size=10,
                 color=money_color(amt), align=PP_ALIGN.RIGHT, fill=stripe, bold=True)
        set_cell(tbl.cell(ri, 4), "收入" if amt > 0 else "支出",
                 size=10, align=PP_ALIGN.CENTER, fill=stripe,
                 color=money_color(amt))
        set_cell(tbl.cell(ri, 5), remark, size=9, fill=stripe, color=GRAY)
        tbl.rows[ri].height = Inches(0.32)
    # 小计行
    sub_row = len(data) + 1
    # 合并 0..2
    cell0 = tbl.cell(sub_row, 0)
    cell0.merge(tbl.cell(sub_row, 2))
    set_cell(cell0, subtotal_label, bold=True, size=12,
             color=DARK_RED, align=PP_ALIGN.CENTER, fill=GOLD)
    set_cell(tbl.cell(sub_row, 3), money_str(subtotal),
             bold=True, size=12, color=money_color(subtotal),
             align=PP_ALIGN.RIGHT, fill=GOLD)
    # 合并 4..5
    cell4 = tbl.cell(sub_row, 4)
    cell4.merge(tbl.cell(sub_row, 5))
    set_cell(cell4, "", fill=GOLD)
    tbl.rows[sub_row].height = Inches(0.46)
    add_footer(s, idx, slide_count)
    return s

# ---- Slide 2: 收支总览 ----
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_top_bar(s, "一、收支总览", "Overview")
# 四个统计卡 + 合计大卡
def stat_card(left, top, w, h, label, value, color):
    bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
    bg.line.color.rgb = GOLD; bg.line.width = Pt(1)
    # 上色条
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, Inches(0.18))
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    tb = s.shapes.add_textbox(left, top + Inches(0.28), w, h - Inches(0.3))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = FONT_CN; r.font.size = Pt(13); r.font.color.rgb = GRAY
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = money_str(value)
    r2.font.name = FONT_CN; r2.font.size = Pt(24); r2.font.bold = True
    r2.font.color.rgb = color
    p2.space_before = Pt(8)

card_w = Inches(2.9); card_h = Inches(1.55); gap = Inches(0.25)
total_w = card_w * 4 + gap * 3
start_left = (SW - total_w) // 2
top_y = Inches(1.45)
stat_card(start_left + (card_w + gap)*0, top_y, card_w, card_h,
          "外部费用", _sum(EXTERNAL), RED)
stat_card(start_left + (card_w + gap)*1, top_y, card_w, card_h,
          "内部费用", _sum(INTERNAL), RED)
stat_card(start_left + (card_w + gap)*2, top_y, card_w, card_h,
          "现场杂项（净额）", _sum(MISC), GREEN if _sum(MISC) >= 0 else RED)
stat_card(start_left + (card_w + gap)*3, top_y, card_w, card_h,
          "赞助收入", _sum(SPONSOR), GREEN)

# 大合计卡片
big_left = Inches(0.6); big_top = Inches(3.25); big_w = SW - Inches(1.2); big_h = Inches(2.4)
big = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, big_left, big_top, big_w, big_h)
big.fill.solid(); big.fill.fore_color.rgb = DARK_RED
big.line.color.rgb = GOLD; big.line.width = Pt(2)

tb = s.shapes.add_textbox(big_left + Inches(0.5), big_top + Inches(0.25),
                          big_w - Inches(1), big_h - Inches(0.5))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "本次活动净结余"
r.font.name = FONT_CN; r.font.size = Pt(18); r.font.color.rgb = GOLD; r.font.bold = True
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = money_str(GRAND_TOTAL)
r2.font.name = FONT_CN; r2.font.size = Pt(54); r2.font.bold = True; r2.font.color.rgb = WHITE
p2.space_before = Pt(8)
p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
r3 = p3.add_run()
r3.text = (f"= 收入 {money_str(_sum(SPONSOR))}  −  支出 ¥{abs(_sum(EXTERNAL)+_sum(INTERNAL)):,.2f}"
           f"  +  代收代付净额 {money_str(_sum(MISC))}")
r3.font.name = FONT_CN; r3.font.size = Pt(14); r3.font.color.rgb = CREAM
p3.space_before = Pt(10)
p4 = tf.add_paragraph(); p4.alignment = PP_ALIGN.CENTER
r4 = p4.add_run(); r4.text = "（负数表示活动总体亏损，需另行筹措补足）"
r4.font.name = FONT_CN; r4.font.size = Pt(11); r4.font.italic = True; r4.font.color.rgb = GOLD
p4.space_before = Pt(8)

add_footer(s, 2, slide_count)

# ---- Slide 3-6: 四张明细表 ----
add_table_slide(3, "二、外部费用明细", f"对外采购 / 服务 / 媒体 / 餐饮  ·  小计 {money_str(_sum(EXTERNAL))}",
                EXTERNAL, "外部费用小计", _sum(EXTERNAL))
add_table_slide(4, "三、内部费用明细", f"内部团队垫付  ·  小计 {money_str(_sum(INTERNAL))}",
                INTERNAL, "内部费用小计", _sum(INTERNAL))
add_table_slide(5, "四、现场杂项（代收代付）", f"升舱代收 + 现场零星支出 + 退费  ·  净额 {money_str(_sum(MISC))}",
                MISC, "现场杂项小计", _sum(MISC))
add_table_slide(6, "五、赞助 / 收入明细", f"企业赞助 + 个人赞助 + 晚宴费  ·  小计 {money_str(_sum(SPONSOR))}",
                SPONSOR, "赞助收入小计", _sum(SPONSOR))

# ---- Slide 7: 证书打印费明细 ----
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_top_bar(s, "附表 · 证书制作及KT版打印明细", "供应商：广告公司  ·  原始单据编号 XH20260522016  ·  合计 ¥1,030.00")
headers = ["序号", "单据编号", "商品名称", "型号", "单位", "数量", "单价", "金额", "备注"]
rows = len(CERT_DETAIL) + 2
cols = len(headers)
left = Inches(0.4); top = Inches(1.4); width = SW - Inches(0.8); height = Inches(4.2)
table_shape = s.shapes.add_table(rows, cols, left, top, width, height)
tbl = table_shape.table
col_w = [0.6, 1.7, 1.4, 0.8, 0.8, 0.8, 1.1, 1.3, 4.0]
total = sum(col_w)
for i, w in enumerate(col_w):
    tbl.columns[i].width = Emu(int((width / total) * w))
for i, h in enumerate(headers):
    set_cell(tbl.cell(0, i), h, bold=True, size=12, color=WHITE, fill=DARK_RED, align=PP_ALIGN.CENTER)
tbl.rows[0].height = Inches(0.42)
for ri, (no, name, mdl, unit, qty, price, amt, remark) in enumerate(CERT_DETAIL, 1):
    stripe = CREAM if ri % 2 == 0 else WHITE
    cells = [ri, no, name, mdl, unit, qty, money_str(price), money_str(amt), remark]
    aligns = [PP_ALIGN.CENTER]*6 + [PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.LEFT]
    for ci, v in enumerate(cells):
        set_cell(tbl.cell(ri, ci), v, size=10, fill=stripe, align=aligns[ci],
                 color=DARK_TXT, bold=(ci==7))
    tbl.rows[ri].height = Inches(0.36)
# 合计行
sr = len(CERT_DETAIL) + 1
c0 = tbl.cell(sr, 0); c0.merge(tbl.cell(sr, 6))
set_cell(c0, "合计", bold=True, size=12, color=DARK_RED, align=PP_ALIGN.CENTER, fill=GOLD)
set_cell(tbl.cell(sr, 7), money_str(1030.00), bold=True, size=12, color=RED, align=PP_ALIGN.RIGHT, fill=GOLD)
set_cell(tbl.cell(sr, 8), "", fill=GOLD)
tbl.rows[sr].height = Inches(0.42)

# 备注块
note_left = Inches(0.4); note_top = Inches(5.85); note_w = SW - Inches(0.8); note_h = Inches(1.15)
nb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, note_left, note_top, note_w, note_h)
nb.fill.solid(); nb.fill.fore_color.rgb = CREAM; nb.line.color.rgb = GOLD; nb.line.width = Pt(0.75)
tb = s.shapes.add_textbox(note_left + Inches(0.2), note_top + Inches(0.1),
                          note_w - Inches(0.4), note_h - Inches(0.2))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "备注："
r.font.name = FONT_CN; r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = DARK_RED
for line in [
    "01-02  证书+外壳：企业 10 件 + 个人 14 件",
    "03     志愿者证书 珠光纸 1 份",
    "04     证书160克超白 + 证书外壳 1 套；闪送：尚九一滴水 1F",
]:
    pn = tf.add_paragraph()
    rn = pn.add_run(); rn.text = "    " + line
    rn.font.name = FONT_CN; rn.font.size = Pt(10); rn.font.color.rgb = DARK_TXT
add_footer(s, 7, slide_count)

# ---- Slide 8: 注意事项 / 结论 ----
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_top_bar(s, "六、注意事项与结论", "Notes & Conclusion")

# 左侧：注意事项
left_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.3),
                               Inches(7.6), Inches(5.6))
left_card.fill.solid(); left_card.fill.fore_color.rgb = WHITE
left_card.line.color.rgb = GOLD; left_card.line.width = Pt(1)
tb = s.shapes.add_textbox(Inches(0.7), Inches(1.45), Inches(7.0), Inches(5.4))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "复核要点"
r.font.name = FONT_CN; r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = DARK_RED

notes_pp = [
    ("01", "数据来源", "源自活动原始流水PDF，逐行核对，各分项小计与总合计均与原表一致。"),
    ("02", "分类口径", "按资金性质分四类：外部费用、内部费用、现场杂项（代收代付）、赞助 / 收入。"),
    ("03", "符号约定", "红色为支出（负数），绿色为收入（正数）。"),
    ("04", "代收代付", "「现场杂项」性质为代收代付：收升舱费 ¥1,440 抵充胸花/对讲机/送货/退费等，最终净额 +¥514.12。"),
    ("05", "媒体支出", "「媒体支出 ¥11,500」中 ¥500 由胡继刚代付，已并入该笔金额。"),
    ("06", "待复核项", "「活动前看场地用餐 +¥148.85」按原票据登记为正数，请确认是否为退款 / 录入差异。"),
]
for no, h, body in notes_pp:
    p = tf.add_paragraph(); p.space_before = Pt(8)
    r = p.add_run(); r.text = f"  {no}  "
    r.font.name = FONT_CN; r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = GOLD
    r2 = p.add_run(); r2.text = h + "  "
    r2.font.name = FONT_CN; r2.font.size = Pt(13); r2.font.bold = True; r2.font.color.rgb = DARK_RED
    r3 = p.add_run(); r3.text = body
    r3.font.name = FONT_CN; r3.font.size = Pt(12); r3.font.color.rgb = DARK_TXT

# 右侧：结论
right_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(1.3),
                                Inches(4.7), Inches(5.6))
right_card.fill.solid(); right_card.fill.fore_color.rgb = DARK_RED
right_card.line.color.rgb = GOLD; right_card.line.width = Pt(1.5)
tb = s.shapes.add_textbox(Inches(8.4), Inches(1.55), Inches(4.3), Inches(5.2))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "结  论"
r.font.name = FONT_CN; r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = GOLD

rows = [
    ("总收入",  money_str(_sum(SPONSOR)), GREEN),
    ("总支出",  f"-¥{abs(_sum(EXTERNAL)+_sum(INTERNAL)):,.2f}", RGBColor(0xFF,0xB6,0xB6)),
    ("代收代付净额", money_str(_sum(MISC)), GREEN if _sum(MISC) >= 0 else RGBColor(0xFF,0xB6,0xB6)),
]
for label, val, color in rows:
    pn = tf.add_paragraph(); pn.alignment = PP_ALIGN.CENTER; pn.space_before = Pt(14)
    rn = pn.add_run(); rn.text = f"{label}\n"
    rn.font.name = FONT_CN; rn.font.size = Pt(12); rn.font.color.rgb = CREAM
    rv = pn.add_run(); rv.text = val
    rv.font.name = FONT_CN; rv.font.size = Pt(20); rv.font.bold = True; rv.font.color.rgb = color

# 分隔线
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.6), Inches(5.55), Inches(3.9), Emu(20000))
line.fill.solid(); line.fill.fore_color.rgb = GOLD; line.line.fill.background()

pn = tf.add_paragraph(); pn.alignment = PP_ALIGN.CENTER; pn.space_before = Pt(22)
rn = pn.add_run(); rn.text = "净  结  余"
rn.font.name = FONT_CN; rn.font.size = Pt(14); rn.font.color.rgb = GOLD
pn = tf.add_paragraph(); pn.alignment = PP_ALIGN.CENTER; pn.space_before = Pt(4)
rn = pn.add_run(); rn.text = money_str(GRAND_TOTAL)
rn.font.name = FONT_CN; rn.font.size = Pt(32); rn.font.bold = True; rn.font.color.rgb = WHITE
pn = tf.add_paragraph(); pn.alignment = PP_ALIGN.CENTER; pn.space_before = Pt(2)
rn = pn.add_run(); rn.text = "（亏损，需另行筹措）"
rn.font.name = FONT_CN; rn.font.size = Pt(11); rn.font.italic = True; rn.font.color.rgb = CREAM

add_footer(s, 8, slide_count)

pptx_path = os.path.join(OUT_DIR, "2026年活动财务整理汇报.pptx")
prs.save(pptx_path)
print("PPT 已生成：", pptx_path)
print("OK")
