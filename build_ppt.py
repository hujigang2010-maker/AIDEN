# -*- coding: utf-8 -*-
"""
生成《钢铁 × 地产 × 产业园区 跨产业撮合与游学转化平台》商业模型 PPT。

运行: python3 build_ppt.py
输出: 跨产业撮合平台_商业模型.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------
# 主题配色（钢铁 / 工业风：深海军蓝 + 钢灰 + 鎏金）
# ----------------------------------------------------------------------------
NAVY      = RGBColor(0x0F, 0x1B, 0x2D)   # 主背景深蓝
NAVY_2    = RGBColor(0x16, 0x2A, 0x45)   # 次级深蓝
STEEL     = RGBColor(0x2E, 0x4A, 0x6E)   # 钢蓝
STEEL_LT  = RGBColor(0x5B, 0x7A, 0xA6)   # 浅钢蓝
GOLD      = RGBColor(0xD9, 0xA8, 0x2B)   # 鎏金（强调）
GOLD_LT   = RGBColor(0xF0, 0xCB, 0x6B)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE  = RGBColor(0xEE, 0xF1, 0xF6)
GREY      = RGBColor(0x9A, 0xA7, 0xB8)
LIGHT_BG  = RGBColor(0xF4, 0xF6, 0xF9)
CARD_BG   = RGBColor(0xFF, 0xFF, 0xFF)
DARKTXT   = RGBColor(0x1B, 0x2A, 0x3A)
RED       = RGBColor(0xC0, 0x4A, 0x3A)

FONT = "微软雅黑"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ----------------------------------------------------------------------------
# 辅助函数
# ----------------------------------------------------------------------------
def _set_font(run, size, color, bold=False, font=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font
    # 中文字体兼容
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', font)


def add_slide(bg=NAVY):
    slide = prs.slides.add_slide(BLANK)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg
    rect.line.fill.background()
    rect.shadow.inherit = False
    return slide


def add_rect(slide, x, y, w, h, color, line_color=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line_color is not None:
        sp.line.color.rgb = line_color
        sp.line.width = Pt(line_w or 1)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def add_text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    """lines: list of (text, size, color, bold) or (text, size, color, bold, space_after)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    for i, spec in enumerate(lines):
        text, size, color, bold = spec[0], spec[1], spec[2], spec[3]
        space_after = spec[4] if len(spec) > 4 else 6
        line_spacing = spec[5] if len(spec) > 5 else None
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        _set_font(run, size, color, bold)
    return tb


def section_header(slide, no, title, subtitle=None, on_dark=True):
    """页眉：编号块 + 标题"""
    tcol = WHITE if on_dark else DARKTXT
    add_rect(slide, Inches(0.6), Inches(0.55), Inches(0.14), Inches(0.62), GOLD)
    chip = add_rect(slide, Inches(0.85), Inches(0.5), Inches(0.7), Inches(0.7), STEEL if on_dark else NAVY)
    tf = chip.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = no; _set_font(r, 26, GOLD, True)
    add_text(slide, Inches(1.7), Inches(0.45), Inches(10.5), Inches(0.85),
             [(title, 30, tcol, True)], anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(1.72), Inches(1.18), Inches(10.5), Inches(0.4),
                 [(subtitle, 13, GOLD if on_dark else STEEL, False)])


def footer(slide, on_dark=True):
    c = GREY if on_dark else STEEL_LT
    add_text(slide, Inches(0.6), Inches(7.0), Inches(8), Inches(0.35),
             [("钢铁 × 地产 × 产业园区 · 跨产业撮合与游学转化平台", 9, c, False)])
    add_text(slide, Inches(11.4), Inches(7.0), Inches(1.5), Inches(0.35),
             [("模型稿 · 控盘版本", 9, c, False)], align=PP_ALIGN.RIGHT)


# ============================================================================
# 1. 封面
# ============================================================================
s = add_slide(NAVY)
# 渐变质感：叠加深蓝斜块
tri = add_rect(s, Inches(7.6), 0, Inches(5.73), SH, NAVY_2, shape=MSO_SHAPE.PARALLELOGRAM)
tri.adjustments[0] = 0.5
add_rect(s, 0, Inches(2.55), SW, Inches(0.05), GOLD)
add_rect(s, Inches(0.9), Inches(1.05), Inches(0.18), Inches(1.2), GOLD)

add_text(s, Inches(1.25), Inches(0.95), Inches(11), Inches(0.6),
         [("跨产业撮合与游学转化平台", 18, GOLD_LT, False)])
add_text(s, Inches(1.2), Inches(2.75), Inches(11.3), Inches(2.2),
         [("钢铁 × 地产 × 产业园区", 50, WHITE, True, 8),
          ("跨产业撮合与游学转化平台", 40, WHITE, True, 8),
          ("—— 可收费的产业服务系统（控盘版模型稿）", 18, GREY, False)])

# 三件控盘标签
labels = ["需求入口", "撮合规则", "收费权"]
for i, t in enumerate(labels):
    x = Inches(1.2 + i * 2.55)
    chip = add_rect(s, x, Inches(5.4), Inches(2.3), Inches(0.62), STEEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    chip.adjustments[0] = 0.5
    tf = chip.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "● " + t; _set_font(r, 15, GOLD_LT, True)

add_text(s, Inches(1.25), Inches(6.55), Inches(11), Inches(0.5),
         [("复旦大学住房政策研究中心 · 杨浦区科技企业联合会  主导  |  中钢协流通分会  协同", 12, GREY, False)])


# ============================================================================
# 2. 项目定位
# ============================================================================
s = add_slide(NAVY)
section_header(s, "01", "项目定位", "以“产业需求驱动”为核心的跨产业对接平台")

add_text(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(1.0),
         [("构建一个以“产业需求驱动”为核心的跨产业对接平台，", 22, WHITE, True, 4),
          ("形成“考察 + 对接 + 项目转化”的闭环机制。", 22, GOLD_LT, True)])

# 三方连接图
boxes = [
    ("钢铁产业链企业", "供给端", STEEL),
    ("房地产 / 园区 / 制造业", "需求端", STEEL),
    ("考察 + 对接 + 项目转化", "闭环机制", GOLD),
]
bx = Inches(0.9); by = Inches(3.4); bw = Inches(3.55); bh = Inches(1.9); gap = Inches(0.55)
for i, (t, sub, col) in enumerate(boxes):
    x = Emu(int(bx) + i * (int(bw) + int(gap)))
    card = add_rect(s, x, by, bw, bh, NAVY_2, line_color=col, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    card.adjustments[0] = 0.06
    add_rect(s, x, by, bw, Inches(0.12), col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, x, Emu(int(by) + Inches(0.45)), bw, Inches(0.5),
             [(sub, 14, GOLD if col == GOLD else STEEL_LT, True)], align=PP_ALIGN.CENTER)
    add_text(s, x, Emu(int(by) + Inches(0.95)), bw, Inches(0.8),
             [(t, 21, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < 2:
        cx = Emu(int(x) + int(bw) + int(gap) // 2 - Inches(0.18))
        ar = add_rect(s, cx, Emu(int(by) + int(bh)//2 - Inches(0.18)), Inches(0.4), Inches(0.4), GOLD, shape=MSO_SHAPE.CHEVRON)
        tf = ar.text_frame; p = tf.paragraphs[0]
        r = p.add_run(); r.text = "×"; _set_font(r, 14, NAVY, True)

add_text(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(0.8),
         [("连接供给端与需求端，把分散的产业资源沉淀为可运营、可收费的对接通道。", 15, GREY, False)])
footer(s)


# ============================================================================
# 3. 核心价值结构（商业闭环）
# ============================================================================
s = add_slide(LIGHT_BG)
section_header(s, "02", "核心价值结构 · 商业闭环", "流量入口 → 核心变现 → 收益核心", on_dark=False)

mods = [
    ("A", "产业游学模块", "流量入口", [
        "标杆企业参访（钢铁/制造/地产/园区）",
        "小规模高质量闭门交流（10–15人）",
        "产业趋势 / 政策 / 应用场景解读"],
     "获取高质量企业决策人", STEEL),
    ("B", "产业撮合模块", "核心变现", [
        "钢铁企业 → 地产/基建/园区采购需求",
        "房地产/园区 → 材料/供应链/技术需求",
        "一对一闭门撮合会议"],
     "形成交易机会", NAVY),
    ("C", "项目转化模块", "收益核心", [
        "重点项目推进（采购/合作开发/试点）",
        "可落地合作跟踪",
        "项目制服务支持"],
     "产生可收费结果", GOLD),
]
cw = Inches(3.95); ch = Inches(4.5); cy = Inches(1.9); cx0 = Inches(0.7); cgap = Inches(0.22)
for i, (tag, title, role, items, effect, col) in enumerate(mods):
    x = Emu(int(cx0) + i * (int(cw) + int(cgap)))
    card = add_rect(s, x, cy, cw, ch, CARD_BG, line_color=RGBColor(0xDD,0xE3,0xEC), line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    card.adjustments[0] = 0.04
    add_rect(s, x, cy, cw, Inches(0.9), col, shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
    badge = add_rect(s, Emu(int(x)+Inches(0.3)), Emu(int(cy)+Inches(0.22)), Inches(0.6), Inches(0.6), WHITE, shape=MSO_SHAPE.OVAL)
    tf = badge.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = tag; _set_font(r, 22, col if col!=GOLD else NAVY, True)
    add_text(s, Emu(int(x)+Inches(1.05)), Emu(int(cy)+Inches(0.16)), Inches(2.7), Inches(0.7),
             [(title, 18, WHITE, True, 0), (role, 12, GOLD_LT if col!=GOLD else NAVY, True)])
    yy = int(cy) + int(Inches(1.25))
    for it in items:
        add_text(s, Emu(int(x)+Inches(0.3)), Emu(yy), Inches(3.45), Inches(0.7),
                 [("• " + it, 13, DARKTXT, False)])
        yy += int(Inches(0.78))
    # 作用条
    eff = add_rect(s, Emu(int(x)+Inches(0.25)), Emu(int(cy)+int(ch)-Inches(0.95)), Emu(int(cw)-int(Inches(0.5))), Inches(0.7),
                   LIGHT_BG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Emu(int(x)+Inches(0.25)), Emu(int(cy)+int(ch)-Inches(0.95)), Emu(int(cw)-int(Inches(0.5))), Inches(0.7),
             [("作用 · " + effect, 14, col if col!=GOLD else RGBColor(0xB0,0x86,0x10), True)],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
footer(s, on_dark=False)


# ============================================================================
# 4. 组织结构（控盘）
# ============================================================================
s = add_slide(NAVY)
section_header(s, "03", "组织结构 · 必须你控盘", "主导方设计规则与收费，渠道方提供资源")

# 主导方（控盘）
lead = add_rect(s, Inches(0.7), Inches(1.75), Inches(5.7), Inches(2.4), NAVY_2, line_color=GOLD, line_w=2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(0.95), Inches(1.9), Inches(5.3), Inches(0.5), [("主导方（控盘方）", 18, GOLD, True)])
add_text(s, Inches(0.95), Inches(2.45), Inches(5.3), Inches(0.9),
         [("复旦大学住房政策研究中心  ·  品牌 + 规则制定", 13, WHITE, False, 4),
          ("杨浦区科技企业联合会  ·  政府与产业资源", 13, WHITE, False)])
add_text(s, Inches(0.95), Inches(3.35), Inches(5.3), Inches(0.7),
         [("职责：设计规则 · 控制内容 · 控制收费 · 控制对外口径", 12, GOLD_LT, True, 0, 1.1)])

# 协同方
coop = add_rect(s, Inches(6.95), Inches(1.75), Inches(5.7), Inches(1.15), NAVY_2, line_color=STEEL_LT, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(7.2), Inches(1.85), Inches(5.2), Inches(0.45), [("协同方（渠道方）", 16, STEEL_LT, True)])
add_text(s, Inches(7.2), Inches(2.3), Inches(5.3), Inches(0.55),
         [("中钢协流通分会 · 提供钢铁企业资源 / 组织参与 / 协助邀约", 12, WHITE, False, 0, 1.1)])

# 执行支持
exe = add_rect(s, Inches(6.95), Inches(3.0), Inches(5.7), Inches(1.15), NAVY_2, line_color=STEEL, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(7.2), Inches(3.1), Inches(5.2), Inches(0.45), [("执行支持", 16, STEEL_LT, True)])
add_text(s, Inches(7.2), Inches(3.55), Inches(5.3), Inches(0.5),
         [("项目运营团队 · 行程 / 会务 / 内容执行", 12, WHITE, False)])

# 控盘三件事
add_text(s, Inches(0.7), Inches(4.55), Inches(11), Inches(0.5),
         [("控盘核心：谁掌握这三件事，谁就掌控全局", 16, WHITE, True)])
three = [("需求入口", "掌握供需两端的进出口"), ("撮合规则", "定义对接与成交的规则"), ("收费权", "占据交易中的收费位置")]
for i, (t, d) in enumerate(three):
    x = Inches(0.7 + i * 4.05)
    card = add_rect(s, x, Inches(5.2), Inches(3.8), Inches(1.35), STEEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, x, Inches(5.32), Inches(3.8), Inches(0.5), [(t, 18, GOLD_LT, True)], align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(5.85), Inches(3.8), Inches(0.5), [(d, 12, OFFWHITE, False)], align=PP_ALIGN.CENTER)
footer(s)


# ============================================================================
# 5. 运行机制（三段式闭环）
# ============================================================================
s = add_slide(LIGHT_BG)
section_header(s, "04", "运行机制 · 三段式闭环", "需求收集 → 线下撮合 → 项目跟踪", on_dark=False)

steps = [
    ("1", "需求收集", ["钢铁企业需求（产品/渠道/场景）", "房地产/园区需求（材料/合作/采购）"], STEEL),
    ("2", "线下撮合", ["游学 + 闭门会议", "1v1 对接机制"], NAVY),
    ("3", "项目跟踪", ["建立项目池", "转入服务或撮合收费"], GOLD),
]
sw = Inches(3.7); sh_ = Inches(2.9); sy = Inches(2.3); sx0 = Inches(0.7); sgap = Inches(0.65)
for i, (no, title, items, col) in enumerate(steps):
    x = Emu(int(sx0) + i * (int(sw) + int(sgap)))
    card = add_rect(s, x, sy, sw, sh_, CARD_BG, line_color=RGBColor(0xDD,0xE3,0xEC), line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    card.adjustments[0] = 0.05
    circ = add_rect(s, Emu(int(x)+int(sw)//2-int(Inches(0.55))), Emu(int(sy)-Inches(0.0)+Inches(0.35)), Inches(1.1), Inches(1.1), col, shape=MSO_SHAPE.OVAL)
    tf = circ.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = no; _set_font(r, 34, WHITE if col!=GOLD else NAVY, True)
    add_text(s, x, Emu(int(sy)+Inches(1.7)), sw, Inches(0.5), [(title, 20, DARKTXT, True)], align=PP_ALIGN.CENTER)
    yy = int(sy) + int(Inches(2.25))
    for it in items:
        add_text(s, Emu(int(x)+Inches(0.3)), Emu(yy), Emu(int(sw)-int(Inches(0.6))), Inches(0.5),
                 [("• " + it, 12.5, RGBColor(0x40,0x50,0x60), False)], align=PP_ALIGN.CENTER)
        yy += int(Inches(0.42))
    if i < 2:
        cx = Emu(int(x) + int(sw) + int(sgap)//2 - int(Inches(0.22)))
        ar = add_rect(s, cx, Emu(int(sy)+int(sh_)//2-int(Inches(0.22))), Inches(0.45), Inches(0.45), GOLD, shape=MSO_SHAPE.CHEVRON)

# 闭环回流
add_rect(s, Inches(0.7), Inches(5.55), Inches(11.9), Inches(0.06), STEEL_LT)
add_text(s, Inches(0.7), Inches(5.75), Inches(11.9), Inches(0.6),
         [("项目池持续回流：成交线索沉淀为可复用资产，反哺下一轮游学与撮合，形成可持续闭环。", 14, STEEL, True)],
         align=PP_ALIGN.CENTER)
footer(s, on_dark=False)


# ============================================================================
# 6. 对外报价逻辑 · 总览
# ============================================================================
s = add_slide(NAVY)
section_header(s, "05", "对外报价逻辑 · 核心武器", "你的收费不是“活动费”，而是“产业撮合服务费”")

quote = [
    ("01", "基础游学套餐", "入口层", "9,800 – 16,800 元/人", STEEL),
    ("02", "产业撮合服务费", "核心利润", "5万–30万/项目 或 1%–3%", GOLD),
    ("03", "企业年度会员", "稳定现金流", "10万 / 30万 每年", STEEL),
    ("04", "高阶专项项目", "利润放大器", "咨询费 + 成交分成", STEEL),
]
qw = Inches(2.85); qh = Inches(3.4); qy = Inches(2.05); qx0 = Inches(0.7); qgap = Inches(0.25)
for i, (no, title, role, price, col) in enumerate(quote):
    x = Emu(int(qx0) + i * (int(qw) + int(qgap)))
    border = GOLD if col == GOLD else STEEL
    card = add_rect(s, x, qy, qw, qh, NAVY_2, line_color=border, line_w=2 if col==GOLD else 1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    card.adjustments[0] = 0.05
    add_text(s, Emu(int(x)+Inches(0.25)), Emu(int(qy)+Inches(0.2)), qw, Inches(0.6), [(no, 26, border, True)])
    add_text(s, Emu(int(x)+Inches(0.25)), Emu(int(qy)+Inches(0.95)), Emu(int(qw)-int(Inches(0.4))), Inches(0.6),
             [(title, 18, WHITE, True)])
    add_text(s, Emu(int(x)+Inches(0.25)), Emu(int(qy)+Inches(1.55)), qw, Inches(0.4), [(role, 13, GOLD_LT, True)])
    pl = add_rect(s, Emu(int(x)+Inches(0.25)), Emu(int(qy)+Inches(2.25)), Emu(int(qw)-int(Inches(0.5))), Inches(0.9), STEEL if col!=GOLD else GOLD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Emu(int(x)+Inches(0.25)), Emu(int(qy)+Inches(2.25)), Emu(int(qw)-int(Inches(0.5))), Inches(0.9),
             [(price, 14, WHITE if col!=GOLD else NAVY, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.7), Inches(5.7), Inches(11.9), Inches(0.6),
         [("逻辑主线：游学筛客 → 撮合变现 → 会员锁定现金流 → 专项项目放大利润。", 15, GREY, False)])
footer(s)


# ============================================================================
# 7. 报价1：基础游学套餐
# ============================================================================
s = add_slide(LIGHT_BG)
section_header(s, "06", "基础游学套餐 · 入口层", "《产业标杆游学（2.5天）》", on_dark=False)

# 左：套餐内容
left = add_rect(s, Inches(0.7), Inches(1.9), Inches(5.9), Inches(4.5), CARD_BG, line_color=RGBColor(0xDD,0xE3,0xEC), line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_rect(s, Inches(0.7), Inches(1.9), Inches(5.9), Inches(0.7), STEEL, shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
add_text(s, Inches(0.95), Inches(2.0), Inches(5.4), Inches(0.5), [("套餐内容", 18, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(1.0), Inches(2.9), Inches(5.4), Inches(3.2),
         [("标准规模：10–15 人", 16, DARKTXT, True, 12),
          ("• 标杆企业参访（钢铁/制造/地产/园区）", 14, RGBColor(0x40,0x50,0x60), False, 10),
          ("• 主题交流：产业趋势 / 政策 / 场景解读", 14, RGBColor(0x40,0x50,0x60), False, 10),
          ("• 闭门对接：高质量决策人面对面", 14, RGBColor(0x40,0x50,0x60), False, 10),
          ("• 周期：2.5 天高密度行程", 14, RGBColor(0x40,0x50,0x60), False, 10)])

# 右：分层定价
right = add_rect(s, Inches(6.85), Inches(1.9), Inches(5.75), Inches(4.5), CARD_BG, line_color=RGBColor(0xDD,0xE3,0xEC), line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_rect(s, Inches(6.85), Inches(1.9), Inches(5.75), Inches(0.7), GOLD, shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
add_text(s, Inches(7.1), Inches(2.0), Inches(5.3), Inches(0.5), [("分层报名费（筛选机制）", 18, NAVY, True)], anchor=MSO_ANCHOR.MIDDLE)
tiers = [("标准票", "9,800 元", STEEL_LT), ("进阶票", "12,800 元", STEEL), ("尊享票", "16,800 元", NAVY)]
for i, (t, p, col) in enumerate(tiers):
    y = Inches(2.85 + i * 1.05)
    bar = add_rect(s, Inches(7.1), y, Inches(5.25), Inches(0.85), col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(7.35), y, Inches(2.5), Inches(0.85), [(t, 16, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(9.3), y, Inches(2.85), Inches(0.85), [(p, 20, GOLD_LT, True)], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(7.1), Inches(6.0), Inches(5.3), Inches(0.4),
         [("用途：覆盖组织成本 + 筛掉低质量客户", 12, STEEL, True)])
footer(s, on_dark=False)


# ============================================================================
# 8. 报价2：产业撮合服务费
# ============================================================================
s = add_slide(NAVY)
section_header(s, "07", "产业撮合服务费 · 核心利润", "只要产生交易，就必须有你的位置")

add_text(s, Inches(0.7), Inches(1.7), Inches(11.8), Inches(0.5),
         [("适用：钢铁企业 vs 房地产/园区采购 · 制造企业 vs 场景合作", 15, GOLD_LT, True)])

# 方案A
a = add_rect(s, Inches(0.7), Inches(2.4), Inches(5.85), Inches(3.4), NAVY_2, line_color=STEEL_LT, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(1.0), Inches(2.65), Inches(5.3), Inches(0.5), [("方案 A · 项目制", 22, WHITE, True)])
add_text(s, Inches(1.0), Inches(3.25), Inches(5.3), Inches(0.5), [("适合单笔、可界定的撮合项目", 13, GREY, False)])
pa = add_rect(s, Inches(1.0), Inches(3.95), Inches(5.25), Inches(1.0), STEEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(1.0), Inches(3.95), Inches(5.25), Inches(1.0),
         [("5 万 – 30 万 / 项目", 26, GOLD_LT, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(1.0), Inches(5.1), Inches(5.3), Inches(0.5), [("✓ 现金流清晰，启动门槛低", 13, OFFWHITE, False)])

# 方案B
b = add_rect(s, Inches(6.75), Inches(2.4), Inches(5.85), Inches(3.4), NAVY_2, line_color=GOLD, line_w=2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(7.05), Inches(2.65), Inches(5.3), Inches(0.5), [("方案 B · 比例制（更高级）", 22, GOLD, True)])
add_text(s, Inches(7.05), Inches(3.25), Inches(5.3), Inches(0.5), [("与成交金额绑定，上不封顶", 13, GREY, False)])
pb = add_rect(s, Inches(7.05), Inches(3.95), Inches(5.25), Inches(1.0), GOLD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(7.05), Inches(3.95), Inches(5.25), Inches(1.0),
         [("成交金额的 1% – 3%", 26, NAVY, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(7.05), Inches(5.1), Inches(5.3), Inches(0.5), [("✓ 绑定大额交易，利润弹性最大", 13, OFFWHITE, False)])

add_text(s, Inches(0.7), Inches(6.1), Inches(11.9), Inches(0.6),
         [("核心逻辑：把“撮合位置”写进规则——你是交易的必经节点，而非可选环节。", 15, GOLD_LT, True)],
         align=PP_ALIGN.CENTER)
footer(s)


# ============================================================================
# 9. 报价3：企业年度会员
# ============================================================================
s = add_slide(LIGHT_BG)
section_header(s, "08", "产业协同会员体系 · 稳定现金流", "用年度机制把客户锁进系统", on_dark=False)

members = [
    ("基础会员", "10 万 / 年", STEEL, [
        "参与闭门会", "常规资源对接", "游学优先报名"]),
    ("高级会员", "30 万 / 年", NAVY, [
        "优先撮合机会", "定向资源对接", "深度参与闭门会", "专项项目优先权"]),
]
mw = Inches(5.7); mh = Inches(4.3); my = Inches(2.0); mx0 = Inches(0.85); mgap = Inches(0.6)
for i, (title, price, col, rights) in enumerate(members):
    x = Emu(int(mx0) + i * (int(mw) + int(mgap)))
    card = add_rect(s, x, my, mw, mh, CARD_BG, line_color=(GOLD if i==1 else RGBColor(0xDD,0xE3,0xEC)), line_w=(2 if i==1 else 1), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_rect(s, x, my, mw, Inches(1.25), col, shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
    add_text(s, Emu(int(x)+Inches(0.4)), Emu(int(my)+Inches(0.18)), Emu(int(mw)-int(Inches(0.8))), Inches(0.5), [(title, 22, WHITE, True)])
    add_text(s, Emu(int(x)+Inches(0.4)), Emu(int(my)+Inches(0.62)), Emu(int(mw)-int(Inches(0.8))), Inches(0.6), [(price, 28, GOLD_LT, True)])
    yy = int(my) + int(Inches(1.6))
    for rgt in rights:
        add_text(s, Emu(int(x)+Inches(0.5)), Emu(yy), Emu(int(mw)-int(Inches(0.9))), Inches(0.5),
                 [("✓  " + rgt, 15, RGBColor(0x35,0x45,0x58), False)])
        yy += int(Inches(0.62))
footer(s, on_dark=False)


# ============================================================================
# 10. 报价4：高阶专项项目
# ============================================================================
s = add_slide(NAVY)
section_header(s, "09", "专项产业项目 · 利润放大器", "把撮合升级为可落地的产业项目")

cases = [
    ("钢铁材料\n地产项目试点", "新材料进入地产应用场景", STEEL),
    ("园区绿色\n材料改造", "存量园区的低碳升级改造", STEEL),
    ("智能制造\n导入地产场景", "制造能力对接地产需求", STEEL),
]
cw2 = Inches(3.85); cy2 = Inches(1.95); cx2 = Inches(0.7); cg2 = Inches(0.28); chh = Inches(2.55)
for i, (t, d, col) in enumerate(cases):
    x = Emu(int(cx2) + i * (int(cw2) + int(cg2)))
    card = add_rect(s, x, cy2, cw2, chh, NAVY_2, line_color=STEEL_LT, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_rect(s, x, cy2, Inches(0.14), chh, GOLD)
    add_text(s, Emu(int(x)+Inches(0.35)), Emu(int(cy2)+Inches(0.4)), Emu(int(cw2)-int(Inches(0.6))), Inches(1.3),
             [(t.replace("\n", " "), 19, WHITE, True, 0, 1.1)])
    add_text(s, Emu(int(x)+Inches(0.35)), Emu(int(cy2)+Inches(1.75)), Emu(int(cw2)-int(Inches(0.6))), Inches(0.6),
             [(d, 13, GREY, False)])

price = add_rect(s, Inches(0.7), Inches(4.95), Inches(11.9), Inches(1.3), STEEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(1.1), Inches(5.0), Inches(11), Inches(0.5), [("收费结构", 15, GOLD_LT, True)], anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(1.1), Inches(5.45), Inches(11), Inches(0.7),
         [("项目咨询费  +  成交分成", 28, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
footer(s)


# ============================================================================
# 11. 分工收益结构
# ============================================================================
s = add_slide(LIGHT_BG)
section_header(s, "10", "分工收益结构 · 你必须掌控", "收费权 + 撮合费 + 项目费 始终在主办方", on_dark=False)

rows = [
    ("你（主办方）", "收费权 + 撮合费 + 项目费", "100% 控盘", GOLD),
    ("中钢协分会", "渠道分成", "20% – 30%", STEEL),
    ("企业参与方", "获得资源对接", "—", STEEL_LT),
]
# 表头
hx = Inches(0.85); hy = Inches(2.1); tw = Inches(11.6)
add_rect(s, hx, hy, tw, Inches(0.7), NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
heads = [("角色", Inches(0.4), Inches(3.2)), ("收益内容", Inches(3.8), Inches(5.0)), ("分成 / 占比", Inches(9.0), Inches(3.2))]
for h, ox, wd in heads:
    add_text(s, Emu(int(hx)+int(ox)), hy, wd, Inches(0.7), [(h, 15, GOLD_LT, True)], anchor=MSO_ANCHOR.MIDDLE)
for i, (role, income, share, col) in enumerate(rows):
    y = Emu(int(hy) + int(Inches(0.85)) + i * int(Inches(1.15)))
    rc = CARD_BG if i % 2 == 0 else RGBColor(0xEC,0xF0,0xF6)
    add_rect(s, hx, y, tw, Inches(1.0), rc, line_color=RGBColor(0xDD,0xE3,0xEC), line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_rect(s, hx, y, Inches(0.12), Inches(1.0), col)
    add_text(s, Emu(int(hx)+int(Inches(0.4))), y, Inches(3.2), Inches(1.0), [(role, 17, DARKTXT, True)], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Emu(int(hx)+int(Inches(3.8))), y, Inches(5.0), Inches(1.0), [(income, 15, RGBColor(0x40,0x50,0x60), False)], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Emu(int(hx)+int(Inches(9.0))), y, Inches(2.4), Inches(1.0), [(share, 18, col if col!=STEEL_LT else STEEL, True)], anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(0.85), Inches(6.55), Inches(11.6), Inches(0.5),
         [("底线：渠道方拿渠道分成，但需求入口、撮合规则与收费权不外放。", 14, STEEL, True)], align=PP_ALIGN.CENTER)
footer(s, on_dark=False)


# ============================================================================
# 12. 核心一句话 / 控盘
# ============================================================================
s = add_slide(NAVY)
add_rect(s, 0, Inches(2.3), SW, Inches(0.05), GOLD)
add_rect(s, 0, Inches(5.0), SW, Inches(0.05), GOLD)
add_text(s, Inches(1.0), Inches(0.9), Inches(11.3), Inches(0.6),
         [("对外统一口径", 16, GOLD, True)], align=PP_ALIGN.CENTER)
add_text(s, Inches(1.0), Inches(2.55), Inches(11.3), Inches(2.3),
         [("我们不是做游学活动，", 32, WHITE, True, 10),
          ("而是做“钢铁产业与地产/园区之间的", 32, WHITE, True, 10),
          ("需求撮合与项目转化平台”。", 34, GOLD_LT, True)], align=PP_ALIGN.CENTER)
add_text(s, Inches(1.0), Inches(5.4), Inches(11.3), Inches(0.6),
         [("本质：建立一个“产业交易入口 + 收费规则系统”", 18, GREY, True)], align=PP_ALIGN.CENTER)
three = ["需求入口（你）", "撮合规则（你）", "收费权（你）"]
for i, t in enumerate(three):
    x = Inches(1.6 + i * 3.5)
    chip = add_rect(s, x, Inches(6.2), Inches(3.1), Inches(0.7), STEEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, x, Inches(6.2), Inches(3.1), Inches(0.7), [("● " + t, 15, GOLD_LT, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============================================================================
# 13. 下一步
# ============================================================================
s = add_slide(LIGHT_BG)
section_header(s, "11", "下一步 · 锁死主动权", "把规则与收费权彻底握在手里", on_dark=False)

nexts = [
    ("对中钢协分会谈判话术", "明确合作边界与对外口径，避免资源被反向拿走"),
    ("分成底线模型", "设定渠道分成上限（20%–30%）与不可让渡条款"),
    ("需求入口锁定", "需求登记与撮合流程统一由主办方归集与分发"),
    ("收费规则固化", "游学 / 撮合 / 会员 / 项目四类收费写入标准协议"),
]
for i, (t, d) in enumerate(nexts):
    col = i % 2; row = i // 2
    x = Inches(0.85 + col * 6.05); y = Inches(2.1 + row * 2.15)
    card = add_rect(s, x, y, Inches(5.7), Inches(1.85), CARD_BG, line_color=RGBColor(0xDD,0xE3,0xEC), line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    num = add_rect(s, Emu(int(x)+Inches(0.3)), Emu(int(y)+Inches(0.35)), Inches(0.9), Inches(0.9), STEEL, shape=MSO_SHAPE.OVAL)
    tf = num.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(i+1); _set_font(r, 26, GOLD_LT, True)
    add_text(s, Emu(int(x)+Inches(1.45)), Emu(int(y)+Inches(0.3)), Inches(4.0), Inches(0.6), [(t, 17, DARKTXT, True)])
    add_text(s, Emu(int(x)+Inches(1.45)), Emu(int(y)+Inches(0.9)), Inches(4.05), Inches(0.8), [(d, 12.5, RGBColor(0x45,0x55,0x66), False, 0, 1.05)])
footer(s, on_dark=False)


out = "跨产业撮合平台_商业模型.pptx"
prs.save(out)
print("Saved:", out, "slides:", len(prs.slides._sldIdLst))
