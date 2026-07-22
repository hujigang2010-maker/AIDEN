# -*- coding: utf-8 -*-
"""为 10 场活动各生成一份独立 PPT（单场立项用）。"""
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import plan_data as D

PLUM = RGBColor(0x2E, 0x1F, 0x47)
PURPLE = RGBColor(0x5B, 0x3E, 0x8E)
PURPLE2 = RGBColor(0x76, 0x58, 0xA8)
LILAC = RGBColor(0xC0, 0xAE, 0xE0)
LAV = RGBColor(0xEC, 0xE6, 0xF7)
LAVED = RGBColor(0xDD, 0xD2, 0xEF)
BGT = RGBColor(0xFB, 0xFA, 0xFE)
GOLD = RGBColor(0xC1, 0x9A, 0x3A)
GREY = RGBColor(0x60, 0x5A, 0x6B)
DARK = RGBColor(0x2B, 0x25, 0x36)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TIER_C = {"A": RGBColor(0x86, 0x6F, 0xB4), "B": PURPLE}

FONT = "Microsoft YaHei"
OUT_DIR = "单场立项PPT"


def _set_font(run, size, bold, color, font=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", font)


def rect(s, x, y, w, h, fill=None, line=None, line_w=None, shadow=False, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w or 1)
    sp.shadow.inherit = False
    if shadow:
        el = sp._element.spPr
        ef = el.makeelement(qn("a:effectLst"), {})
        sh = ef.makeelement(
            qn("a:outerShdw"),
            {"blurRad": "55000", "dist": "25000", "dir": "5400000", "rotWithShape": "0"},
        )
        clr = sh.makeelement(qn("a:srgbClr"), {"val": "7A6BA0"})
        alpha = clr.makeelement(qn("a:alpha"), {"val": "38000"})
        clr.append(alpha)
        sh.append(clr)
        ef.append(sh)
        el.append(ef)
    return sp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for txt, size, bold, color in para:
            r = p.add_run()
            r.text = txt
            _set_font(r, size, bold, color)
    return tb


def bg(s):
    rect(s, 0, 0, SW, SH, fill=BGT)


def header(s, kicker, title):
    rect(s, 0, 0, SW, Inches(1.15), fill=PLUM)
    rect(s, 0, Inches(1.15), SW, Pt(3), fill=GOLD)
    rect(s, Inches(0.55), Inches(0.48), Pt(4), Inches(0.48), fill=GOLD)
    text(s, Inches(0.72), Inches(0.16), Inches(12), Inches(0.3),
         [[(kicker, 11, True, LILAC)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(0.72), Inches(0.44), Inches(12), Inches(0.6),
         [[(title, 22, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)


def footer(s, a, n, total=5):
    rect(s, Inches(0.55), Inches(7.06), Inches(12.25), Pt(0.75), fill=LAVED)
    text(s, Inches(0.55), Inches(7.1), Inches(10.5), Inches(0.3),
         [[(f"{a['oa_code']}  |  {a['title']}  |  单场立项", 8, False, GREY)]])
    text(s, Inches(11.8), Inches(7.1), Inches(1.0), Inches(0.3),
         [[(f"{n}/{total}", 8, True, PURPLE)]], align=PP_ALIGN.RIGHT)


def safe_name(a):
    t = re.sub(r'[\\/:*?"<>|]', "", a["title"])
    t = t.replace("·", "-").replace("“", "").replace("”", "").replace("（", "(").replace("）", ")")
    return f"{a['no']:02d}_{a['oa_code']}_{a['month']}_{t}"


def build_one(a):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    tc = a["tier"].split(" ")[0]
    accent = TIER_C[tc]

    # ---------- 1 封面 ----------
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, SW, SH, fill=PLUM)
    rect(s, 0, 0, SW, Inches(2.5), fill=PURPLE)
    rect(s, 0, Inches(2.5), SW, Pt(4), fill=GOLD)
    text(s, Inches(0.72), Inches(0.5), Inches(12), Inches(0.4),
         [[(f"单场活动立项方案  ·  {a['oa_code']}", 14, True, LILAC)]])
    text(s, Inches(0.72), Inches(1.05), Inches(12), Inches(1.2),
         [[(f"NO.{a['no']}  {a['title']}", 30, True, WHITE)]], space_after=2)
    rect(s, Inches(0.75), Inches(2.85), Inches(1.8), Pt(3), fill=GOLD)
    text(s, Inches(0.72), Inches(3.1), Inches(12), Inches(0.45),
         [[("复旦大学  ·  上海市科技企业联合会  ·  东方枢纽", 14, True, WHITE)]])
    meta_line = f"{a['date']}  |  {a['scale']}  |  {a['sector']}  |  报价 {a['price']} 万元"
    text(s, Inches(0.72), Inches(3.7), Inches(12), Inches(0.4),
         [[(meta_line, 13, False, LILAC)]])
    text(s, Inches(0.72), Inches(4.4), Inches(12), Inches(0.4),
         [[(f"场地：{a['venue']}", 12, False, LILAC)]])
    text(s, Inches(0.72), Inches(6.5), Inches(12), Inches(0.4),
         [[("东方枢纽 A 片区整体办公招商  |  单场独立走 OA", 11, False, LILAC)]])

    # ---------- 2 基本信息 / OA 要件 ----------
    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, f"单场立项 · {a['oa_code']}", "基本信息与 OA 要件")
    # 左信息卡
    rect(s, Inches(0.55), Inches(1.45), Inches(6.0), Inches(5.2), fill=WHITE, line=LAVED, line_w=1, shadow=True)
    rect(s, Inches(0.55), Inches(1.45), Inches(6.0), Inches(0.55), fill=accent)
    text(s, Inches(0.55), Inches(1.45), Inches(6.0), Inches(0.55),
         [[("活动基本信息", 14, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rows = [
        ("立项编号", a["oa_code"]),
        ("活动主题", a["title"]),
        ("拟定时间", a["date"]),
        ("所属月份", a["month"]),
        ("产业板块", a["sector"]),
        ("形式规模", a["scale"]),
        ("场地建议", a["venue"]),
        ("报价档位", a["tier"].split(" · ")[0] + " 档"),
        ("单场报价", f"{a['price']} 万元"),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.15) + i * Inches(0.45)
        fill = LAV if i % 2 == 0 else WHITE
        rect(s, Inches(0.7), y, Inches(5.7), Inches(0.42), fill=fill)
        text(s, Inches(0.85), y, Inches(1.6), Inches(0.42), [[(k, 11, True, GREY)]], anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(2.5), y, Inches(3.7), Inches(0.42),
             [[(v, 11, True if i == 8 else False, GOLD if i == 8 else DARK)]],
             anchor=MSO_ANCHOR.MIDDLE)

    # 右 OA 勾选
    rect(s, Inches(6.85), Inches(1.45), Inches(5.95), Inches(5.2), fill=WHITE, line=LAVED, line_w=1, shadow=True)
    rect(s, Inches(6.85), Inches(1.45), Inches(5.95), Inches(0.55), fill=PLUM)
    text(s, Inches(6.85), Inches(1.45), Inches(5.95), Inches(0.55),
         [[("OA 提交要件（勾选）", 14, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    checks = [
        "活动主题与时间",
        "参会规模与场地",
        "邀约名单（企业 / 嘉宾）",
        "参会人员背景（决策层）",
        "单场预算明细（见第 4 页）",
        "指定策划供应商签约路径",
    ]
    for i, c in enumerate(checks):
        y = Inches(2.25) + i * Inches(0.65)
        rect(s, Inches(7.1), y, Inches(5.45), Inches(0.55), fill=LAV)
        rect(s, Inches(7.1), y, Pt(5), Inches(0.55), fill=GOLD)
        text(s, Inches(7.35), y, Inches(5.0), Inches(0.55),
             [[(f"☐  {c}", 13, True, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
    footer(s, a, 2)

    # ---------- 3 内容策划 ----------
    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, f"单场立项 · {a['oa_code']}", "内容策划 · 嘉宾资源 · 招商衔接")
    # 内容
    rect(s, Inches(0.55), Inches(1.45), Inches(12.25), Inches(1.85), fill=WHITE, line=LAVED, line_w=1, shadow=True)
    rect(s, Inches(0.55), Inches(1.45), Inches(12.25), Inches(0.45), fill=accent)
    text(s, Inches(0.75), Inches(1.45), Inches(12), Inches(0.45),
         [[("内容建议", 13, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    runs = [[("· ", 12, True, GOLD), (ln, 12, False, DARK)] for ln in a["content"]]
    text(s, Inches(0.8), Inches(2.05), Inches(11.8), Inches(1.1), runs, space_after=4, line_spacing=1.15)

    # 嘉宾 + 价值 两列
    rect(s, Inches(0.55), Inches(3.5), Inches(6.0), Inches(3.15), fill=WHITE, line=LAVED, line_w=1, shadow=True)
    rect(s, Inches(0.55), Inches(3.5), Inches(6.0), Inches(0.45), fill=PURPLE)
    text(s, Inches(0.75), Inches(3.5), Inches(5.6), Inches(0.45),
         [[("拟邀嘉宾 / 资源", 13, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    runs = [[("· ", 11, True, GOLD), (g, 11, False, DARK)] for g in a["guests"]]
    text(s, Inches(0.8), Inches(4.15), Inches(5.5), Inches(2.3), runs, space_after=6, line_spacing=1.15)

    rect(s, Inches(6.8), Inches(3.5), Inches(6.0), Inches(3.15), fill=WHITE, line=LAVED, line_w=1, shadow=True)
    rect(s, Inches(6.8), Inches(3.5), Inches(6.0), Inches(0.45), fill=PLUM)
    text(s, Inches(7.0), Inches(3.5), Inches(5.6), Inches(0.45),
         [[("招商衔接价值", 13, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    runs = [[("· ", 11, True, GOLD), (ln, 11, False, DARK)] for ln in a["invest"]]
    text(s, Inches(7.05), Inches(4.15), Inches(5.5), Inches(1.5), runs, space_after=6, line_spacing=1.15)
    text(s, Inches(7.05), Inches(5.8), Inches(5.5), Inches(0.65),
         [[("立项说明：", 11, True, PURPLE), (a["oa_note"], 11, False, GREY)]], line_spacing=1.1)
    footer(s, a, 3)

    # ---------- 4 预算明细 ----------
    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, f"单场立项 · {a['oa_code']}", "单场预算明细")
    text(s, Inches(0.72), Inches(1.35), Inches(12), Inches(0.35),
         [[(f"档位：{a['tier']}    |    单位：万元    |    最终以指定策划供应商合同为准", 11, False, GREY)]])

    # 表头
    y0 = Inches(1.85)
    rect(s, Inches(1.5), y0, Inches(7.0), Inches(0.5), fill=PLUM)
    rect(s, Inches(8.5), y0, Inches(3.3), Inches(0.5), fill=PLUM)
    text(s, Inches(1.5), y0, Inches(7.0), Inches(0.5),
         [[("费用构成项", 13, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(8.5), y0, Inches(3.3), Inches(0.5),
         [[("金额（万元）", 13, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    for i, (item, val) in enumerate(zip(D.COST_ITEMS, a["costs"])):
        y = y0 + Inches(0.5) * (i + 1)
        fill = LAV if i % 2 == 0 else WHITE
        rect(s, Inches(1.5), y, Inches(7.0), Inches(0.5), fill=fill, line=LAVED, line_w=0.5)
        rect(s, Inches(8.5), y, Inches(3.3), Inches(0.5), fill=fill, line=LAVED, line_w=0.5)
        text(s, Inches(1.7), y, Inches(6.6), Inches(0.5), [[(item, 12, False, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(8.5), y, Inches(3.3), Inches(0.5),
             [[(str(val), 13, True, GOLD)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    yt = y0 + Inches(0.5) * (len(D.COST_ITEMS) + 1)
    rect(s, Inches(1.5), yt, Inches(7.0), Inches(0.55), fill=accent)
    rect(s, Inches(8.5), yt, Inches(3.3), Inches(0.55), fill=GOLD)
    text(s, Inches(1.5), yt, Inches(7.0), Inches(0.55),
         [[("单场合计", 14, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(8.5), yt, Inches(3.3), Inches(0.55),
         [[(str(a["price"]), 16, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, a, 4)

    # ---------- 5 执行与下一步 ----------
    s = prs.slides.add_slide(blank)
    bg(s)
    header(s, f"单场立项 · {a['oa_code']}", "执行安排与下一步")
    steps = [
        ("① 确认档期", f"敲定 {a['date']} 具体日期与场地档期"),
        ("② 提交 OA", "提交邀约名单、参会人数、人员背景、本场预算明细"),
        ("③ 供应商签约", "经东方枢纽指定策划供应商签约支付，落实搭建通行证"),
        ("④ 精准邀约", "按决策层标准完成企业/嘉宾邀约与确认"),
        ("⑤ 现场执行", "会务执行、签到注册、数据全量录入后台"),
        ("⑥ 会后跟进", "导入 CRM，完成意向跟进与会后评估"),
    ]
    for i, (t, d) in enumerate(steps):
        col, row = i % 2, i // 2
        x = Inches(0.55) + col * Inches(6.35)
        y = Inches(1.5) + row * Inches(1.55)
        w, h = Inches(6.05), Inches(1.35)
        rect(s, x, y, w, h, fill=WHITE, line=LAVED, line_w=1, shadow=True)
        rect(s, x, y, Inches(0.12), h, fill=accent)
        text(s, x + Inches(0.35), y + Inches(0.25), w - Inches(0.5), Inches(0.4),
             [[(t, 14, True, PURPLE)]])
        text(s, x + Inches(0.35), y + Inches(0.7), w - Inches(0.55), Inches(0.5),
             [[(d, 12, False, DARK)]])
    footer(s, a, 5)

    os.makedirs(OUT_DIR, exist_ok=True)
    path = os.path.join(OUT_DIR, safe_name(a) + ".pptx")
    prs.save(path)
    return path


# 全局尺寸（在函数外供 footer/header 使用）
SW = Inches(13.333)
SH = Inches(7.5)

if __name__ == "__main__":
    paths = []
    for a in D.ACTIVITIES:
        p = build_one(a)
        paths.append(p)
        print("已生成:", p)
    print(f"\n合计 {len(paths)} 份单场 PPT → 目录：{OUT_DIR}/")
