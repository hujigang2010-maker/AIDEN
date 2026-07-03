# -*- coding: utf-8 -*-
"""生成《项目概要》PPT(16:9)。"""
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

import content as C

OUT_DIR = os.path.join(os.path.dirname(__file__), "..", "deliverables")
CN_FONT = "微软雅黑"
ACCENT = RGBColor(0x1F, 0x4E, 0x79)
GRAY = RGBColor(0x59, 0x59, 0x59)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xD9, 0xE2, 0xF3)

SLIDE_W, SLIDE_H = Cm(33.87), Cm(19.05)


def add_textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def set_run(r, text, size, bold=False, color=None):
    r.text = text
    r.font.name = CN_FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color if color else RGBColor(0x26, 0x26, 0x26)


def add_para(tf, text, size=16, bold=False, color=None, bullet_char=None,
             space_after=8, align=None, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    r = p.add_run()
    set_run(r, (f"{bullet_char} {text}" if bullet_char else text), size, bold, color)
    p.space_after = Pt(space_after)
    if align:
        p.alignment = align
    return p


def add_bar(slide, top=Cm(0), height=Cm(0.35), color=ACCENT):
    from pptx.enum.shapes import MSO_SHAPE
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), top, SLIDE_W, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def content_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bar(slide)
    tf = add_textbox(slide, Cm(1.5), Cm(0.8), Cm(30.8), Cm(1.8))
    add_para(tf, title, size=28, bold=True, color=ACCENT, first=True)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ---------------- 封面 ----------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    from pptx.enum.shapes import MSO_SHAPE
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = ACCENT
    bg.line.fill.background()
    tf = add_textbox(slide, Cm(2.5), Cm(5.2), Cm(28.8), Cm(8))
    add_para(tf, C.PROJECT_NAME, size=34, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, first=True, space_after=16)
    add_para(tf, "项目概要", size=48, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, space_after=28)
    add_para(tf, f"提供方:{C.PROVIDER_LINE}", size=18, color=LIGHT,
             align=PP_ALIGN.CENTER, space_after=6)
    add_para(tf, C.DOC_DATE, size=16, color=LIGHT, align=PP_ALIGN.CENTER)

    # ---------------- 概念 ----------------
    slide = content_slide(prs, f"一、{C.CONCEPT_TITLE}")
    tf = add_textbox(slide, Cm(1.8), Cm(3.2), Cm(30.3), Cm(14.5))
    for pt in C.CONCEPT_POINTS:
        add_para(tf, pt, size=19, bullet_char="■", space_after=18,
                 first=(pt is C.CONCEPT_POINTS[0]))

    # ---------------- 产业方向 ----------------
    slide = content_slide(prs, f"二、{C.INDUSTRY_TITLE}")
    tf = add_textbox(slide, Cm(1.8), Cm(3.0), Cm(30.3), Cm(1.4))
    add_para(tf, C.INDUSTRY_INTRO, size=18, color=GRAY, first=True)
    tf = add_textbox(slide, Cm(1.8), Cm(4.6), Cm(30.3), Cm(13.0))
    for name, detail in C.INDUSTRY_GROUPS:
        p = tf.paragraphs[0] if (name, detail) == C.INDUSTRY_GROUPS[0] and \
            not tf.paragraphs[0].runs else tf.add_paragraph()
        r = p.add_run()
        set_run(r, f"■ {name}:", 20, bold=True, color=ACCENT)
        r2 = p.add_run()
        set_run(r2, detail, 19)
        p.space_after = Pt(18)

    # ---------------- 下阶段重点 ----------------
    slide = content_slide(prs, f"三、{C.NEXT_STAGE_TITLE}")
    tf = add_textbox(slide, Cm(1.8), Cm(3.2), Cm(30.3), Cm(2.6))
    add_para(tf, C.NEXT_STAGE_TEXT, size=20, first=True)
    tf = add_textbox(slide, Cm(1.8), Cm(6.2), Cm(30.3), Cm(1.2))
    add_para(tf, "策划服务八大模块", size=20, bold=True, color=ACCENT, first=True)
    # 两列展示八大模块
    half = (len(C.SERVICE_MODULES) + 1) // 2
    for col, items in enumerate((C.SERVICE_MODULES[:half], C.SERVICE_MODULES[half:])):
        tf = add_textbox(slide, Cm(1.8) + Cm(15.5) * col, Cm(7.6), Cm(15.0), Cm(10.5))
        for i, (name, _) in enumerate(items):
            idx = col * half + i + 1
            add_para(tf, f"{idx}. {name}", size=17, space_after=12, first=(i == 0))

    # ---------------- 成果交付 ----------------
    slide = content_slide(prs, "四、成果交付")
    tf = add_textbox(slide, Cm(1.8), Cm(3.2), Cm(30.3), Cm(14.0))
    for i, d in enumerate(C.DELIVERABLES):
        add_para(tf, d, size=19, bullet_char="■", space_after=18, first=(i == 0))

    # ---------------- 提资清单摘要 ----------------
    slide = content_slide(prs, "五、提资清单(摘要)")
    tf = add_textbox(slide, Cm(1.8), Cm(3.0), Cm(30.3), Cm(1.2))
    add_para(tf, "完整清单详见 Excel《提资清单与报价》。高优先级资料如下:",
             size=16, color=GRAY, first=True)
    high_items = [(c_, n) for c_, n, _, p_ in C.INFO_REQUEST_ITEMS if p_ == "高"]
    half = (len(high_items) + 1) // 2
    for col, items in enumerate((high_items[:half], high_items[half:])):
        tf = add_textbox(slide, Cm(1.8) + Cm(15.5) * col, Cm(4.4), Cm(15.0), Cm(13.5))
        for i, (cat, name) in enumerate(items):
            p = tf.paragraphs[0] if i == 0 and not tf.paragraphs[0].runs \
                else tf.add_paragraph()
            r = p.add_run()
            set_run(r, f"■ [{cat}] ", 16, bold=True, color=ACCENT)
            r2 = p.add_run()
            set_run(r2, name, 16)
            p.space_after = Pt(14)

    # ---------------- 报价 ----------------
    slide = content_slide(prs, "六、服务报价")
    rows, cols = len(C.QUOTATION_ITEMS) + 2, 3
    table_shape = slide.shapes.add_table(
        rows, cols, Cm(1.8), Cm(3.2), Cm(30.3), Cm(9.5))
    table = table_shape.table
    table.columns[0].width = Cm(9.5)
    table.columns[1].width = Cm(15.3)
    table.columns[2].width = Cm(5.5)
    headers = ["工作阶段", "主要工作内容", "报价(万元)"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        set_run(cell.text_frame.paragraphs[0].add_run(), h, 15, bold=True, color=WHITE)
    for i, (stage, work, _, fee) in enumerate(C.QUOTATION_ITEMS, 1):
        for j, v in enumerate([stage, work, str(fee)]):
            cell = table.cell(i, j)
            cell.text = ""
            set_run(cell.text_frame.paragraphs[0].add_run(), v, 13)
    cell = table.cell(rows - 1, 0)
    cell.text = ""
    set_run(cell.text_frame.paragraphs[0].add_run(), "合计", 15, bold=True)
    cell = table.cell(rows - 1, 2)
    cell.text = ""
    set_run(cell.text_frame.paragraphs[0].add_run(),
            f"{C.QUOTATION_TOTAL}", 15, bold=True, color=ACCENT)
    tf = add_textbox(slide, Cm(1.8), Cm(13.2), Cm(30.3), Cm(5.0))
    for i, n in enumerate(C.QUOTATION_NOTES):
        add_para(tf, n, size=13, color=GRAY, bullet_char="•",
                 space_after=6, first=(i == 0))

    # ---------------- 结尾 ----------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = ACCENT
    bg.line.fill.background()
    tf = add_textbox(slide, Cm(2.5), Cm(7.0), Cm(28.8), Cm(6))
    add_para(tf, "期待与您的合作", size=40, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, first=True, space_after=20)
    add_para(tf, C.PROVIDER_LINE, size=18, color=LIGHT, align=PP_ALIGN.CENTER)

    os.makedirs(OUT_DIR, exist_ok=True)
    out = os.path.join(OUT_DIR, "广州黄埔九佛TOD全球自贸365街区_项目概要.pptx")
    prs.save(out)
    print("saved:", out)


if __name__ == "__main__":
    main()
