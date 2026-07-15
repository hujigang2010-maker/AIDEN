# -*- coding: utf-8 -*-
"""
基于 build_finance.py 的"仅支出"版本：
- 完全剔除「赞助 / 收入」整类
- 剔除「内部费用」中的正向流水（活动前看场地用餐 +148.85）
- 剔除「现场杂项」中的代收升舱费（+1440）
保留所有支出明细，输出 Excel 与 PPT。
"""
import os
import copy
from build_finance import (
    EXTERNAL, INTERNAL, MISC, CERT_DETAIL,
    Workbook, Font, PatternFill, Alignment, Border, Side, NamedStyle,
    get_column_letter,
    COLOR_PRIMARY, COLOR_ACCENT, COLOR_SUBHEAD, COLOR_STRIPE,
    COLOR_TOTAL_BG, COLOR_INCOME, COLOR_EXPENSE,
    thin, medium, BORDER_CELL, BORDER_HEAD,
    FONT_TITLE, FONT_SUBTITLE, FONT_HEAD, FONT_SUBHEAD, FONT_CELL,
    FONT_TOTAL, FONT_NUM_EXP, FONT_NUM_INC,
    FILL_TITLE, FILL_HEAD, FILL_SUBHEAD, FILL_STRIPE, FILL_TOTAL,
    ALIGN_CENTER, ALIGN_LEFT, ALIGN_RIGHT, NUM_FMT,
)
from openpyxl.styles import PatternFill as _PF

OUT_DIR = "/workspace/output"
os.makedirs(OUT_DIR, exist_ok=True)

# ---------------------------------------------------------------
# 1. 仅保留支出项（amt < 0）
# ---------------------------------------------------------------
def expense_only(rows):
    return [r for r in rows if r[2] < 0]

EXTERNAL_E = expense_only(EXTERNAL)
INTERNAL_E = expense_only(INTERNAL)
MISC_E     = expense_only(MISC)

def _sum(rows):
    return round(sum(r[2] for r in rows), 2)

SUM_EXT  = _sum(EXTERNAL_E)
SUM_INT  = _sum(INTERNAL_E)
SUM_MISC = _sum(MISC_E)
GRAND    = round(SUM_EXT + SUM_INT + SUM_MISC, 2)

# 控制台核对
print(f"外部费用（仅支出）共 {len(EXTERNAL_E)} 项，合计 ¥{SUM_EXT:,.2f}")
print(f"内部费用（仅支出）共 {len(INTERNAL_E)} 项，合计 ¥{SUM_INT:,.2f}")
print(f"现场杂项（仅支出）共 {len(MISC_E)} 项，合计 ¥{SUM_MISC:,.2f}")
print(f"总支出合计 ¥{GRAND:,.2f}")

# ---------------------------------------------------------------
# 2. 构建 Excel
# ---------------------------------------------------------------
wb = Workbook()

def _title_row(ws, title, subtitle, ncols):
    ws.row_dimensions[1].height = 38
    ws.row_dimensions[2].height = 22
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=ncols)
    ws.merge_cells(start_row=2, start_column=1, end_row=2, end_column=ncols)
    c = ws.cell(row=1, column=1, value=title)
    c.font = FONT_TITLE; c.fill = FILL_TITLE; c.alignment = ALIGN_CENTER
    c2 = ws.cell(row=2, column=1, value=subtitle)
    c2.font = FONT_SUBTITLE; c2.alignment = ALIGN_CENTER
    for col in range(1, ncols + 1):
        ws.cell(row=1, column=col).border = BORDER_HEAD

def _header_row(ws, row, headers):
    ws.row_dimensions[row].height = 28
    for i, h in enumerate(headers, 1):
        c = ws.cell(row=row, column=i, value=h)
        c.font = FONT_HEAD; c.fill = FILL_HEAD
        c.alignment = ALIGN_CENTER; c.border = BORDER_HEAD

def _write_money(cell, value):
    cell.number_format = NUM_FMT
    cell.value = value
    if value is None or value == 0:
        cell.font = FONT_CELL
    elif value < 0:
        cell.font = FONT_NUM_EXP
    else:
        cell.font = FONT_NUM_INC
    cell.alignment = ALIGN_RIGHT
    cell.border = BORDER_CELL

def _set_widths(ws, widths):
    for i, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = w

def _enable_fit_to_page(ws):
    ws.page_setup.orientation = ws.ORIENTATION_LANDSCAPE
    ws.page_setup.paperSize = ws.PAPERSIZE_A4
    ws.page_setup.fitToWidth = 1
    ws.page_setup.fitToHeight = 0
    ws.sheet_properties.pageSetUpPr.fitToPage = True
    ws.page_margins.left = 0.3
    ws.page_margins.right = 0.3
    ws.page_margins.top = 0.4
    ws.page_margins.bottom = 0.4
    ws.print_options.horizontalCentered = True

# ---------- 工作表 1：支出总览 ----------
ws = wb.active
ws.title = "支出总览"
_title_row(
    ws,
    "2026 年度活动 · 支出明细汇总（仅支出）",
    "口径：剔除所有进账（赞助 / 代收 / 退款），仅展示对外、对内、现场支出  ｜  币种：人民币（元）",
    5,
)
_header_row(ws, 3, ["序号", "类别", "项目数", "支出金额（元）", "说明"])

summary_rows = [
    ("一", "外部费用",  len(EXTERNAL_E), SUM_EXT,  "对外采购、服务、媒体、餐饮等"),
    ("二", "内部费用",  len(INTERNAL_E), SUM_INT,  "内部团队垫付的物料、伴手礼、餐饮等"),
    ("三", "现场杂项",  len(MISC_E),     SUM_MISC, "胸花/对讲机/退费/送货等现场零星支出"),
]
r = 4
for idx, (no, name, cnt, amt, desc) in enumerate(summary_rows):
    ws.cell(row=r, column=1, value=no).alignment = ALIGN_CENTER
    ws.cell(row=r, column=2, value=name).alignment = ALIGN_LEFT
    ws.cell(row=r, column=3, value=cnt).alignment = ALIGN_CENTER
    _write_money(ws.cell(row=r, column=4), amt)
    ws.cell(row=r, column=5, value=desc).alignment = ALIGN_LEFT
    for col in range(1, 6):
        cell = ws.cell(row=r, column=col)
        cell.border = BORDER_CELL
        if col != 4:
            cell.font = FONT_CELL
        if idx % 2 == 1 and col != 4:
            cell.fill = FILL_STRIPE
    r += 1

# 合计
ws.row_dimensions[r].height = 30
ws.merge_cells(start_row=r, start_column=1, end_row=r, end_column=3)
total_label = ws.cell(row=r, column=1, value="支出合计")
total_label.font = FONT_TOTAL; total_label.fill = FILL_TOTAL
total_label.alignment = ALIGN_CENTER; total_label.border = BORDER_CELL
_write_money(ws.cell(row=r, column=4), GRAND)
ws.cell(row=r, column=4).font = FONT_TOTAL
ws.cell(row=r, column=4).fill = FILL_TOTAL
note = ws.cell(row=r, column=5, value=f"共 {len(EXTERNAL_E)+len(INTERNAL_E)+len(MISC_E)} 笔支出累加（外 + 内 + 现场）")
note.font = FONT_TOTAL; note.fill = FILL_TOTAL
note.alignment = ALIGN_LEFT; note.border = BORDER_CELL

# 结论
r += 2
ws.merge_cells(start_row=r, start_column=1, end_row=r, end_column=5)
c = ws.cell(row=r, column=1,
            value=f"本表为纯支出口径：本次活动总支出 ¥{abs(GRAND):,.2f}（以正额展示为支出，以负号展示为现金流出）。"
                  f"如需含进账的完整收支表请见另一版本《2026年活动财务整理表.xlsx》。")
c.font = Font(name="WenQuanYi Micro Hei", size=11, bold=True, color="8B0000")
c.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)
c.fill = FILL_SUBHEAD
ws.row_dimensions[r].height = 38

_set_widths(ws, [6, 22, 10, 18, 60])
ws.freeze_panes = "A4"
ws.sheet_view.showGridLines = False
_enable_fit_to_page(ws)

# ---------- 通用：明细表 ----------
def build_detail_sheet(sheet_name, title, subtitle, data, subtotal_label, subtotal_value):
    ws = wb.create_sheet(sheet_name)
    _title_row(ws, title, subtitle, 5)
    _header_row(ws, 3, ["序号", "项目", "对接人 / 供应商", "支出金额（元）", "备注"])
    r = 4
    for i, (proj, contact, amt, remark) in enumerate(data, 1):
        ws.row_dimensions[r].height = 22 if len(remark) < 40 else 36
        ws.cell(row=r, column=1, value=i).alignment = ALIGN_CENTER
        ws.cell(row=r, column=2, value=proj).alignment = ALIGN_LEFT
        ws.cell(row=r, column=3, value=contact).alignment = ALIGN_LEFT
        _write_money(ws.cell(row=r, column=4), amt)
        ws.cell(row=r, column=5, value=remark).alignment = ALIGN_LEFT
        for col in range(1, 6):
            cell = ws.cell(row=r, column=col)
            cell.border = BORDER_CELL
            if col != 4:
                cell.font = FONT_CELL
            if i % 2 == 0 and col != 4:
                cell.fill = FILL_STRIPE
        r += 1
    # 小计
    ws.row_dimensions[r].height = 28
    ws.merge_cells(start_row=r, start_column=1, end_row=r, end_column=3)
    sl = ws.cell(row=r, column=1, value=subtotal_label)
    sl.font = FONT_TOTAL; sl.fill = FILL_TOTAL; sl.alignment = ALIGN_CENTER; sl.border = BORDER_CELL
    _write_money(ws.cell(row=r, column=4), subtotal_value)
    ws.cell(row=r, column=4).font = FONT_TOTAL
    ws.cell(row=r, column=4).fill = FILL_TOTAL
    c = ws.cell(row=r, column=5, value="")
    c.fill = FILL_TOTAL; c.border = BORDER_CELL
    _set_widths(ws, [6, 22, 22, 16, 65])
    ws.freeze_panes = "A4"
    ws.sheet_view.showGridLines = False
    _enable_fit_to_page(ws)

build_detail_sheet(
    "外部费用明细", "一、外部费用明细（仅支出）",
    f"对外采购 / 服务 / 媒体 / 餐饮  ·  共 {len(EXTERNAL_E)} 项",
    EXTERNAL_E, "外部费用小计", SUM_EXT,
)
build_detail_sheet(
    "内部费用明细", "二、内部费用明细（仅支出）",
    f"内部团队垫付物料、伴手礼、餐饮等  ·  共 {len(INTERNAL_E)} 项  ｜  已剔除"
    "「活动前看场地用餐 +148.85」",
    INTERNAL_E, "内部费用小计", SUM_INT,
)
build_detail_sheet(
    "现场杂项明细", "三、现场杂项明细（仅支出）",
    f"胸花 / 对讲机 / 退费 / 送货等  ·  共 {len(MISC_E)} 项  ｜  已剔除"
    "「升舱费代收 +1,440」",
    MISC_E, "现场杂项小计", SUM_MISC,
)

# ---------- 工作表：证书打印费明细 ----------
ws = wb.create_sheet("证书打印费明细")
_title_row(
    ws,
    "附表 · 证书制作及KT版打印明细",
    "供应商：广告公司  ｜  原始单据编号：XH20260522016  ｜  合计 ¥1,030.00",
    8,
)
_header_row(ws, 3, ["序号", "单据编号", "商品名称", "型号", "单位", "数量", "单价（元）", "金额（元）"])
r = 4
for i, (no, name, mdl, unit, qty, price, amt, _remark) in enumerate(CERT_DETAIL, 1):
    ws.cell(row=r, column=1, value=i).alignment = ALIGN_CENTER
    ws.cell(row=r, column=2, value=no).alignment = ALIGN_CENTER
    ws.cell(row=r, column=3, value=name).alignment = ALIGN_LEFT
    ws.cell(row=r, column=4, value=mdl).alignment = ALIGN_CENTER
    ws.cell(row=r, column=5, value=unit).alignment = ALIGN_CENTER
    ws.cell(row=r, column=6, value=qty).alignment = ALIGN_CENTER
    _write_money(ws.cell(row=r, column=7), price)
    ws.cell(row=r, column=7).font = FONT_CELL
    _write_money(ws.cell(row=r, column=8), amt)
    ws.cell(row=r, column=8).font = FONT_CELL
    for col in range(1, 9):
        cell = ws.cell(row=r, column=col)
        cell.border = BORDER_CELL
        if cell.font is None:
            cell.font = FONT_CELL
        if i % 2 == 0:
            cell.fill = FILL_STRIPE
    r += 1
ws.merge_cells(start_row=r, start_column=1, end_row=r, end_column=7)
sl = ws.cell(row=r, column=1, value="合计")
sl.font = FONT_TOTAL; sl.fill = FILL_TOTAL; sl.alignment = ALIGN_CENTER; sl.border = BORDER_CELL
_write_money(ws.cell(row=r, column=8), 1030.00)
ws.cell(row=r, column=8).font = FONT_TOTAL
ws.cell(row=r, column=8).fill = FILL_TOTAL
_set_widths(ws, [6, 18, 14, 8, 8, 8, 14, 16])
ws.freeze_panes = "A4"
ws.sheet_view.showGridLines = False
_enable_fit_to_page(ws)

# ---------- 阅读说明 ----------
ws = wb.create_sheet("阅读说明", 0)
_title_row(
    ws,
    "2026 年度活动 · 支出明细汇总表（仅支出版）",
    "本表仅列示活动支出项，已剔除所有进账（赞助、代收、退款等）",
    2,
)
notes = [
    ("整理人 / 日期",   "活动财务组  ·  2026 年"),
    ("币种",           "人民币（元）"),
    ("数据来源",       "活动原始流水 PDF"),
    ("口径说明",       "本表仅保留资金流出项目（支出），剔除所有进账类目"),
    ("剔除明细 1",     "整类剔除「赞助 / 收入」(¥+46,758.00)"),
    ("剔除明细 2",     "「现场杂项」中代收升舱费 ¥+1,440 已剔除（同类的 -199 退费亦在支出口径下保留）"),
    ("剔除明细 3",     "「内部费用」中「活动前看场地用餐 +¥148.85」已剔除"),
    ("分类",           "① 外部费用  ② 内部费用  ③ 现场杂项"),
    ("总支出",         f"¥{abs(GRAND):,.2f}（= ¥{abs(SUM_EXT):,.2f} + ¥{abs(SUM_INT):,.2f} + ¥{abs(SUM_MISC):,.2f}）"),
    ("配套文件",       "完整收支版：2026年活动财务整理表.xlsx"),
    ("工作表索引",     "支出总览 → 外部费用明细 → 内部费用明细 → 现场杂项明细 → 证书打印费明细"),
]
r = 4
for k, v in notes:
    ws.row_dimensions[r].height = 26
    a = ws.cell(row=r, column=1, value=k)
    a.font = FONT_SUBHEAD; a.fill = FILL_SUBHEAD
    a.alignment = ALIGN_CENTER; a.border = BORDER_CELL
    b = ws.cell(row=r, column=2, value=v)
    b.font = FONT_CELL; b.alignment = ALIGN_LEFT; b.border = BORDER_CELL
    r += 1
_set_widths(ws, [18, 72])
ws.sheet_view.showGridLines = False
_enable_fit_to_page(ws)

xlsx_path = os.path.join(OUT_DIR, "2026年活动财务整理表_仅支出.xlsx")
wb.save(xlsx_path)
print("Excel 已生成：", xlsx_path)

# ---------------------------------------------------------------
# 3. PPT 汇报
# ---------------------------------------------------------------
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

DARK_RED = RGBColor(0x8B, 0x00, 0x00)
GOLD     = RGBColor(0xC9, 0xA8, 0x6B)
CREAM    = RGBColor(0xFB, 0xF7, 0xF0)
DARK_TXT = RGBColor(0x33, 0x33, 0x33)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0x6B, 0x6B, 0x6B)
RED      = RGBColor(0xB2, 0x22, 0x22)
BG_PAGE  = RGBColor(0xFB, 0xF7, 0xF0)

FONT_CN = "WenQuanYi Micro Hei"

def add_bg(slide, color=BG_PAGE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg

def add_top_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.95))
    bar.fill.solid(); bar.fill.fore_color.rgb = DARK_RED
    bar.line.fill.background()
    deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.95), SW, Inches(0.06))
    deco.fill.solid(); deco.fill.fore_color.rgb = GOLD
    deco.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), SW - Inches(1), Inches(0.85))
    tf = tb.text_frame; tf.margin_left = 0; tf.margin_top = 0
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.name = FONT_CN; r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.name = FONT_CN; r2.font.size = Pt(12); r2.font.color.rgb = GOLD

def add_footer(slide, idx, total):
    tb = slide.shapes.add_textbox(Inches(0.4), SH - Inches(0.4), SW - Inches(0.8), Inches(0.3))
    tf = tb.text_frame; tf.margin_left = 0; tf.margin_top = 0
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "2026 年度活动 · 支出明细汇总（仅支出版）"
    r.font.name = FONT_CN; r.font.size = Pt(9); r.font.color.rgb = GRAY
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run(); r2.text = f"{idx} / {total}"
    r2.font.name = FONT_CN; r2.font.size = Pt(9); r2.font.color.rgb = GRAY

def set_cell(cell, text, *, bold=False, size=11, color=DARK_TXT,
             fill=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    cell.text = ""
    tf = cell.text_frame
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(text)
    r.font.name = FONT_CN; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    cell.vertical_anchor = anchor
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = fill

def money_str(x):
    return f"¥{x:,.2f}"

slide_count = 6

# ---- Slide 1: 封面 ----
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG_PAGE)
top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(3.6))
top.fill.solid(); top.fill.fore_color.rgb = DARK_RED; top.line.fill.background()
deco = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.6), SW, Inches(0.08))
deco.fill.solid(); deco.fill.fore_color.rgb = GOLD; deco.line.fill.background()

tb = s.shapes.add_textbox(Inches(0.8), Inches(1.0), SW - Inches(1.6), Inches(2.4))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "2026 年度活动"
r.font.name = FONT_CN; r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = GOLD
p2 = tf.add_paragraph()
r2 = p2.add_run(); r2.text = "支 出 明 细 汇 总 报 告"
r2.font.name = FONT_CN; r2.font.size = Pt(54); r2.font.bold = True; r2.font.color.rgb = WHITE
p3 = tf.add_paragraph()
r3 = p3.add_run(); r3.text = "Expense-only Report  ·  All Income Excluded"
r3.font.name = "Calibri"; r3.font.size = Pt(14); r3.font.color.rgb = CREAM

card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.2), SW - Inches(1.6), Inches(2.6))
card.fill.solid(); card.fill.fore_color.rgb = WHITE
card.line.color.rgb = GOLD; card.line.width = Pt(1.5)
tb = s.shapes.add_textbox(Inches(1.2), Inches(4.4), SW - Inches(2.4), Inches(2.3))
tf = tb.text_frame; tf.word_wrap = True
for key, val, c in [
    ("口径说明", "本表仅展示活动支出，已剔除所有进账（赞助 / 代收 / 退款）", DARK_TXT),
    ("外部费用", money_str(SUM_EXT), RED),
    ("内部费用", money_str(SUM_INT), RED),
    ("现场杂项", money_str(SUM_MISC), RED),
    ("总支出",   f"{money_str(GRAND)}  （绝对值 ¥{abs(GRAND):,.2f}）", RED),
]:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    r = p.add_run(); r.text = f"  {key} ： "
    r.font.name = FONT_CN; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = DARK_RED
    r2 = p.add_run(); r2.text = val
    r2.font.name = FONT_CN; r2.font.size = Pt(14); r2.font.color.rgb = c
    p.space_after = Pt(6)

add_footer(s, 1, slide_count)

# ---- Slide 2: 支出总览 ----
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_top_bar(s, "一、支出总览", "Expense Overview · 仅支出口径")

def stat_card(left, top, w, h, label, value, count):
    bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
    bg.line.color.rgb = GOLD; bg.line.width = Pt(1)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, Inches(0.18))
    bar.fill.solid(); bar.fill.fore_color.rgb = RED
    bar.line.fill.background()
    tb = s.shapes.add_textbox(left, top + Inches(0.28), w, h - Inches(0.3))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = FONT_CN; r.font.size = Pt(13); r.font.color.rgb = GRAY
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = money_str(value)
    r2.font.name = FONT_CN; r2.font.size = Pt(24); r2.font.bold = True
    r2.font.color.rgb = RED
    p2.space_before = Pt(8)
    p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run(); r3.text = f"共 {count} 项"
    r3.font.name = FONT_CN; r3.font.size = Pt(11); r3.font.color.rgb = GRAY

card_w = Inches(3.8); card_h = Inches(1.85); gap = Inches(0.3)
total_w = card_w * 3 + gap * 2
start_left = (SW - total_w) // 2
top_y = Inches(1.45)
stat_card(start_left + (card_w + gap)*0, top_y, card_w, card_h,
          "外部费用", SUM_EXT, len(EXTERNAL_E))
stat_card(start_left + (card_w + gap)*1, top_y, card_w, card_h,
          "内部费用", SUM_INT, len(INTERNAL_E))
stat_card(start_left + (card_w + gap)*2, top_y, card_w, card_h,
          "现场杂项", SUM_MISC, len(MISC_E))

# 大合计卡
big_left = Inches(0.6); big_top = Inches(3.55); big_w = SW - Inches(1.2); big_h = Inches(2.4)
big = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, big_left, big_top, big_w, big_h)
big.fill.solid(); big.fill.fore_color.rgb = DARK_RED
big.line.color.rgb = GOLD; big.line.width = Pt(2)

tb = s.shapes.add_textbox(big_left + Inches(0.5), big_top + Inches(0.25),
                          big_w - Inches(1), big_h - Inches(0.5))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "本次活动总支出"
r.font.name = FONT_CN; r.font.size = Pt(18); r.font.color.rgb = GOLD; r.font.bold = True
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = money_str(GRAND)
r2.font.name = FONT_CN; r2.font.size = Pt(54); r2.font.bold = True; r2.font.color.rgb = WHITE
p2.space_before = Pt(8)
p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
r3 = p3.add_run()
r3.text = (f"= 外部 {money_str(SUM_EXT)}  +  内部 {money_str(SUM_INT)}  +  现场 {money_str(SUM_MISC)}"
           f"  ·  共 {len(EXTERNAL_E)+len(INTERNAL_E)+len(MISC_E)} 笔")
r3.font.name = FONT_CN; r3.font.size = Pt(14); r3.font.color.rgb = CREAM
p3.space_before = Pt(10)
p4 = tf.add_paragraph(); p4.alignment = PP_ALIGN.CENTER
r4 = p4.add_run(); r4.text = "（已剔除所有进账：赞助 ¥+46,758、升舱代收 ¥+1,440、看场地 ¥+148.85）"
r4.font.name = FONT_CN; r4.font.size = Pt(11); r4.font.italic = True; r4.font.color.rgb = GOLD
p4.space_before = Pt(8)

add_footer(s, 2, slide_count)

# ---- 复用：支出明细表 ----
def add_table_slide(idx, title, subtitle, data, subtotal_label, subtotal):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_top_bar(s, title, subtitle)
    headers = ["序号", "项目", "对接人 / 供应商", "支出金额（元）", "备注"]
    rows = len(data) + 2
    cols = len(headers)
    left = Inches(0.4); top = Inches(1.25)
    width = SW - Inches(0.8); height = Inches(5.6)
    table_shape = s.shapes.add_table(rows, cols, left, top, width, height)
    tbl = table_shape.table
    col_widths_in = [0.7, 2.5, 2.2, 1.9, 5.2]
    total_in = sum(col_widths_in)
    for i, w_in in enumerate(col_widths_in):
        tbl.columns[i].width = Emu(int((width / total_in) * w_in))
    for i, h in enumerate(headers):
        set_cell(tbl.cell(0, i), h, bold=True, size=12, color=WHITE,
                 fill=DARK_RED, align=PP_ALIGN.CENTER)
    tbl.rows[0].height = Inches(0.42)
    for ri, (proj, contact, amt, remark) in enumerate(data, 1):
        stripe = CREAM if ri % 2 == 0 else WHITE
        set_cell(tbl.cell(ri, 0), ri, size=10, align=PP_ALIGN.CENTER, fill=stripe)
        set_cell(tbl.cell(ri, 1), proj, size=10, fill=stripe)
        set_cell(tbl.cell(ri, 2), contact, size=10, fill=stripe)
        set_cell(tbl.cell(ri, 3), money_str(amt), size=10,
                 color=RED, align=PP_ALIGN.RIGHT, fill=stripe, bold=True)
        set_cell(tbl.cell(ri, 4), remark, size=9, fill=stripe, color=GRAY)
        tbl.rows[ri].height = Inches(0.32)
    sub_row = len(data) + 1
    cell0 = tbl.cell(sub_row, 0)
    cell0.merge(tbl.cell(sub_row, 2))
    set_cell(cell0, subtotal_label, bold=True, size=12,
             color=DARK_RED, align=PP_ALIGN.CENTER, fill=GOLD)
    set_cell(tbl.cell(sub_row, 3), money_str(subtotal),
             bold=True, size=12, color=RED,
             align=PP_ALIGN.RIGHT, fill=GOLD)
    set_cell(tbl.cell(sub_row, 4), "", fill=GOLD)
    tbl.rows[sub_row].height = Inches(0.46)
    add_footer(s, idx, slide_count)

add_table_slide(3, "二、外部费用明细", f"对外采购 / 服务 / 媒体 / 餐饮  ·  小计 {money_str(SUM_EXT)}",
                EXTERNAL_E, "外部费用小计", SUM_EXT)
add_table_slide(4, "三、内部费用明细", f"内部团队垫付  ·  小计 {money_str(SUM_INT)}  ·  已剔除 +148.85",
                INTERNAL_E, "内部费用小计", SUM_INT)
add_table_slide(5, "四、现场杂项明细", f"胸花 / 对讲机 / 退费 / 送货  ·  小计 {money_str(SUM_MISC)}  ·  已剔除 +1,440 代收",
                MISC_E, "现场杂项小计", SUM_MISC)

# ---- Slide 6: 复核要点 ----
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_top_bar(s, "五、复核要点与总额", "Notes & Total")

# 左卡
left_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.3),
                               Inches(7.6), Inches(5.6))
left_card.fill.solid(); left_card.fill.fore_color.rgb = WHITE
left_card.line.color.rgb = GOLD; left_card.line.width = Pt(1)
tb = s.shapes.add_textbox(Inches(0.7), Inches(1.45), Inches(7.0), Inches(5.4))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "复核要点"
r.font.name = FONT_CN; r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = DARK_RED

notes_pp = [
    ("01", "数据来源", "源自活动原始流水 PDF；本版本仅保留资金流出（支出）项目。"),
    ("02", "口径变更", "整类剔除「赞助 / 收入」(¥+46,758.00)。"),
    ("03", "剔除明细", "「现场杂项」中代收升舱费 ¥+1,440 已剔除；「内部费用」中「活动前看场地用餐 +¥148.85」已剔除。"),
    ("04", "保留口径", "支出类的负值条目全部保留，包括退费 ¥-199（视为现金流出口径下的支出）。"),
    ("05", "分类",     "保留三大支出类目：外部费用 / 内部费用 / 现场杂项；附「证书打印费明细」附表。"),
    ("06", "配套文件", "如需含进账的完整收支版，请见《2026年活动财务整理表.xlsx》。"),
]
for no, h, body in notes_pp:
    p = tf.add_paragraph(); p.space_before = Pt(8)
    r = p.add_run(); r.text = f"  {no}  "
    r.font.name = FONT_CN; r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = GOLD
    r2 = p.add_run(); r2.text = h + "  "
    r2.font.name = FONT_CN; r2.font.size = Pt(13); r2.font.bold = True; r2.font.color.rgb = DARK_RED
    r3 = p.add_run(); r3.text = body
    r3.font.name = FONT_CN; r3.font.size = Pt(12); r3.font.color.rgb = DARK_TXT

# 右卡
right_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(1.3),
                                Inches(4.7), Inches(5.6))
right_card.fill.solid(); right_card.fill.fore_color.rgb = DARK_RED
right_card.line.color.rgb = GOLD; right_card.line.width = Pt(1.5)
tb = s.shapes.add_textbox(Inches(8.4), Inches(1.55), Inches(4.3), Inches(5.2))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "总  支  出"
r.font.name = FONT_CN; r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = GOLD

rows = [
    ("外部费用", money_str(SUM_EXT)),
    ("内部费用", money_str(SUM_INT)),
    ("现场杂项", money_str(SUM_MISC)),
]
for label, val in rows:
    pn = tf.add_paragraph(); pn.alignment = PP_ALIGN.CENTER; pn.space_before = Pt(14)
    rn = pn.add_run(); rn.text = f"{label}\n"
    rn.font.name = FONT_CN; rn.font.size = Pt(12); rn.font.color.rgb = CREAM
    rv = pn.add_run(); rv.text = val
    rv.font.name = FONT_CN; rv.font.size = Pt(20); rv.font.bold = True
    rv.font.color.rgb = RGBColor(0xFF, 0xB6, 0xB6)

line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.6), Inches(5.55), Inches(3.9), Emu(20000))
line.fill.solid(); line.fill.fore_color.rgb = GOLD; line.line.fill.background()

pn = tf.add_paragraph(); pn.alignment = PP_ALIGN.CENTER; pn.space_before = Pt(22)
rn = pn.add_run(); rn.text = "合  计"
rn.font.name = FONT_CN; rn.font.size = Pt(14); rn.font.color.rgb = GOLD
pn = tf.add_paragraph(); pn.alignment = PP_ALIGN.CENTER; pn.space_before = Pt(4)
rn = pn.add_run(); rn.text = money_str(GRAND)
rn.font.name = FONT_CN; rn.font.size = Pt(32); rn.font.bold = True; rn.font.color.rgb = WHITE
pn = tf.add_paragraph(); pn.alignment = PP_ALIGN.CENTER; pn.space_before = Pt(2)
rn = pn.add_run(); rn.text = f"（绝对值 ¥{abs(GRAND):,.2f}）"
rn.font.name = FONT_CN; rn.font.size = Pt(11); rn.font.italic = True; rn.font.color.rgb = CREAM

add_footer(s, 6, slide_count)

pptx_path = os.path.join(OUT_DIR, "2026年活动财务整理汇报_仅支出.pptx")
prs.save(pptx_path)
print("PPT 已生成：", pptx_path)
print("OK")
