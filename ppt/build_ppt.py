# -*- coding: utf-8 -*-
"""生成《面向"十五五"的不动产转型与高质量发展：从空间开发向科技赋能平台迈进》汇报 PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ============================================================
# 视觉系统
# ============================================================
NAVY        = RGBColor(0x0F, 0x2A, 0x47)
NAVY_DARK   = RGBColor(0x08, 0x1A, 0x2E)
NAVY_LIGHT  = RGBColor(0x1B, 0x3B, 0x6F)
GOLD        = RGBColor(0xC9, 0xA9, 0x61)
ORANGE      = RGBColor(0xE0, 0x7A, 0x3C)
TEAL        = RGBColor(0x2A, 0x9D, 0x8F)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
PAPER       = RGBColor(0xF7, 0xF6, 0xF1)
BG_LIGHT    = RGBColor(0xF5, 0xF7, 0xFA)
GRAY_LINE   = RGBColor(0xE5, 0xE7, 0xEB)
TEXT_DARK   = RGBColor(0x1F, 0x29, 0x37)
TEXT_BODY   = RGBColor(0x37, 0x41, 0x51)
TEXT_MUTE   = RGBColor(0x6B, 0x72, 0x80)

CN_FONT     = '微软雅黑'
EN_FONT     = 'Calibri'

# 中文成对引号常量，避免在源码字符串中混入 ASCII "
LQ = '\u201c'  # "
RQ = '\u201d'  # "

def q(s):
    """把任意普通字符串里的占位 ^^...^^ 转成中文引号 "..."。"""
    out = []
    flag = False
    i = 0
    while i < len(s):
        if s[i:i+2] == '^^':
            out.append(LQ if not flag else RQ)
            flag = not flag
            i += 2
        else:
            out.append(s[i])
            i += 1
    return ''.join(out)


# ============================================================
# 基础工具
# ============================================================
def set_font(run, *, size=18, bold=False, color=TEXT_DARK, name=CN_FONT):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    eastAsia = rPr.find(qn('a:eastAsia'))
    if eastAsia is None:
        eastAsia = etree.SubElement(rPr, qn('a:eastAsia'))
    eastAsia.set('typeface', name)


def add_textbox(slide, left, top, width, height,
                text, *, size=18, bold=False, color=TEXT_DARK,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                line_spacing=1.25):
    text = q(text)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        set_font(run, size=size, bold=bold, color=color)
    return box


def add_rect(slide, left, top, width, height, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


def slide_blank(prs, bg=WHITE):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, bg)
    return slide


def add_header(slide, prs, section_no, section_title, page_no, total_pages):
    sw = prs.slide_width
    sh = prs.slide_height
    add_rect(slide, Inches(0.6), Inches(0.45), Inches(0.06), Inches(0.32), GOLD)
    add_textbox(slide, Inches(0.78), Inches(0.42), Inches(8), Inches(0.4),
                f'PART 0{section_no}  ·  {section_title}',
                size=11, color=TEXT_MUTE, bold=True)
    add_textbox(slide, Inches(8.2), Inches(0.42), Inches(5.2), Inches(0.4),
                '面向^^十五五^^的不动产转型与高质量发展',
                size=10, color=TEXT_MUTE, align=PP_ALIGN.RIGHT)
    add_rect(slide, Inches(0.6), sh - Inches(0.55),
             sw - Inches(1.2), Emu(9525), GRAY_LINE)
    add_textbox(slide, Inches(0.6), sh - Inches(0.45),
                Inches(8), Inches(0.32),
                '从空间开发  ·  迈向科技赋能平台',
                size=9, color=TEXT_MUTE)
    add_textbox(slide, sw - Inches(2.0), sh - Inches(0.45),
                Inches(1.4), Inches(0.32),
                f'{page_no:02d} / {total_pages:02d}',
                size=10, color=NAVY, bold=True, align=PP_ALIGN.RIGHT)


def add_title(slide, main, sub=None, *, top=Inches(1.0)):
    add_textbox(slide, Inches(0.78), top, Inches(12), Inches(0.7),
                main, size=26, bold=True, color=NAVY)
    if sub:
        add_textbox(slide, Inches(0.78), top + Inches(0.7),
                    Inches(12), Inches(0.4),
                    sub, size=12, color=TEXT_MUTE)
    add_rect(slide, Inches(0.78), top + Inches(1.14),
             Inches(0.6), Emu(28575), GOLD)


# ============================================================
# 演示文稿初始化（16:9）
# ============================================================
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height


# ============================================================
# 1. 封面
# ============================================================
def build_cover():
    slide = slide_blank(prs, bg=NAVY_DARK)
    add_rect(slide, 0, 0, Inches(8.2), SH, NAVY)
    add_rect(slide, Inches(8.2), 0, SW - Inches(8.2), SH, PAPER)

    add_rect(slide, Inches(0.7), Inches(0.7), Inches(1.2), Emu(28575), GOLD)
    add_textbox(slide, Inches(0.7), Inches(0.85), Inches(6), Inches(0.35),
                'REAL  ESTATE  ·  HIGH-QUALITY  DEVELOPMENT',
                size=10, color=GOLD, bold=True)

    add_textbox(slide, Inches(0.7), Inches(1.7), Inches(7.4), Inches(1.6),
                '面向^^十五五^^的不动产转型\n与高质量发展',
                size=40, bold=True, color=WHITE, line_spacing=1.15)

    add_rect(slide, Inches(0.7), Inches(4.05), Inches(0.6), Emu(28575), GOLD)
    add_textbox(slide, Inches(0.7), Inches(4.18), Inches(7.4), Inches(0.7),
                '从空间开发  ·  迈向科技赋能平台',
                size=20, color=GOLD)

    add_textbox(slide, Inches(0.7), Inches(5.0), Inches(7.4), Inches(1.5),
                '存量时代的精细化运营 × 大模型驱动的管理革命\n×  科创生态的产城融合',
                size=14, color=WHITE, line_spacing=1.6)

    add_rect(slide, Inches(8.7), Inches(1.7), Inches(0.06), Inches(2.0), GOLD)
    add_textbox(slide, Inches(8.95), Inches(1.7), Inches(4.0), Inches(0.4),
                '汇报主题', size=10, color=TEXT_MUTE, bold=True)
    add_textbox(slide, Inches(8.95), Inches(2.05), Inches(4.0), Inches(1.6),
                '新模式 · 新动能 · 新格局\n房地产高质量发展的实践路径',
                size=14, bold=True, color=NAVY, line_spacing=1.4)

    add_textbox(slide, Inches(8.95), Inches(4.0), Inches(4.0), Inches(0.4),
                '汇 报 提 纲', size=10, color=TEXT_MUTE, bold=True)
    items = [
        '01  时代定调 · 周期新旧转换',
        '02  破局寻路 · 新模式内涵',
        '03  核心动能 · 科技驱动',
        '04  总结展望 · 知与行',
    ]
    add_textbox(slide, Inches(8.95), Inches(4.35), Inches(4.0), Inches(2.0),
                '\n'.join(items),
                size=12, color=TEXT_BODY, line_spacing=1.7)

    add_rect(slide, Inches(0.7), SH - Inches(0.95), Inches(7.4), Emu(9525), GOLD)
    add_textbox(slide, Inches(0.7), SH - Inches(0.78), Inches(7.4), Inches(0.4),
                '汇报人： ___________     |     2026 · ^^十五五^^开局之年',
                size=11, color=WHITE)


# ============================================================
# 2. 目录
# ============================================================
def build_toc():
    slide = slide_blank(prs, bg=BG_LIGHT)
    add_rect(slide, 0, 0, Inches(4.6), SH, NAVY)
    add_textbox(slide, Inches(0.7), Inches(0.9), Inches(3.5), Inches(0.4),
                'C O N T E N T S', size=12, bold=True, color=GOLD)
    add_textbox(slide, Inches(0.7), Inches(1.4), Inches(3.5), Inches(1.5),
                '汇 报\n提 纲', size=44, bold=True, color=WHITE, line_spacing=1.05)
    add_rect(slide, Inches(0.7), Inches(4.0), Inches(0.5), Emu(28575), GOLD)
    add_textbox(slide, Inches(0.7), Inches(4.15), Inches(3.5), Inches(2.5),
                '^^十五五^^开局之年\n房地产高质量发展\n核心议题汇报',
                size=14, color=WHITE, line_spacing=1.6)

    chapters = [
        ('01', '时代定调', '—— ^^十五五^^开局与房地产周期的新旧转换',
         '宏观政策切入  ·  告别旧模式  ·  提出核心论点'),
        ('02', '破局寻路', '—— 房地产^^新模式^^的核心内涵与高质量标准',
         '空间载体演变  ·  轻重分离  ·  产城融合'),
        ('03', '核心动能', '—— 科技创新如何驱动不动产高质量发展',
         'AI 与大模型  ·  科创生态赋能  ·  跨界协同'),
        ('04', '总结展望', '—— ^^十五五^^时期从业者的知与行',
         '发展共识  ·  行业阵痛  ·  新质生产力土壤'),
    ]
    top = Inches(0.85)
    row_h = Inches(1.45)
    for i, (no, title, sub, tag) in enumerate(chapters):
        y = top + row_h * i
        add_rect(slide, Inches(4.95), y, Inches(1.2), Inches(1.2), WHITE, line=GRAY_LINE)
        add_textbox(slide, Inches(4.95), y + Inches(0.18),
                    Inches(1.2), Inches(0.5),
                    no, size=30, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_rect(slide, Inches(5.25), y + Inches(0.85),
                 Inches(0.6), Emu(19050), GOLD)
        add_textbox(slide, Inches(6.4), y + Inches(0.12),
                    Inches(6.5), Inches(0.45),
                    title, size=20, bold=True, color=NAVY)
        add_textbox(slide, Inches(6.4), y + Inches(0.55),
                    Inches(6.5), Inches(0.4),
                    sub, size=12, color=TEXT_BODY)
        add_textbox(slide, Inches(6.4), y + Inches(0.92),
                    Inches(6.5), Inches(0.4),
                    '▎ ' + tag, size=10, color=GOLD, bold=True)


# ============================================================
# 章节封面
# ============================================================
def build_section_divider(no, title, en, sub, key_points):
    slide = slide_blank(prs, bg=NAVY)
    add_rect(slide, 0, 0, Inches(0.5), SH, GOLD)
    add_textbox(slide, Inches(7.5), Inches(0.4), Inches(5.5), Inches(4.0),
                f'0{no}', size=240, bold=True, color=NAVY_LIGHT,
                align=PP_ALIGN.RIGHT)

    add_textbox(slide, Inches(1.0), Inches(1.5), Inches(2.0), Inches(0.4),
                f'PART  ·  0{no}', size=12, bold=True, color=GOLD)
    add_textbox(slide, Inches(1.0), Inches(1.95), Inches(8), Inches(0.4),
                en, size=11, color=GOLD)

    add_textbox(slide, Inches(1.0), Inches(2.6), Inches(11), Inches(1.2),
                title, size=36, bold=True, color=WHITE, line_spacing=1.15)

    add_rect(slide, Inches(1.0), Inches(4.3), Inches(0.6), Emu(28575), GOLD)
    add_textbox(slide, Inches(1.0), Inches(4.45), Inches(11), Inches(0.5),
                sub, size=15, color=WHITE)

    box_top = Inches(5.3)
    box_h = Inches(1.4)
    box_w = Inches(3.8)
    gap = Inches(0.15)
    for i, kp in enumerate(key_points):
        x = Inches(1.0) + (box_w + gap) * i
        add_rect(slide, x, box_top, box_w, box_h, NAVY_LIGHT, line=None)
        add_rect(slide, x, box_top, Inches(0.06), box_h, GOLD)
        add_textbox(slide, x + Inches(0.25), box_top + Inches(0.18),
                    box_w - Inches(0.4), Inches(0.4),
                    kp[0], size=12, bold=True, color=GOLD)
        add_textbox(slide, x + Inches(0.25), box_top + Inches(0.55),
                    box_w - Inches(0.4), box_h - Inches(0.65),
                    kp[1], size=11, color=WHITE, line_spacing=1.45)


# ============================================================
# 通用内容页
# ============================================================
def new_content_slide(section_no, section_title, page_no, total_pages, title, sub=None):
    slide = slide_blank(prs, bg=WHITE)
    add_header(slide, prs, section_no, section_title, page_no, total_pages)
    add_title(slide, title, sub)
    return slide


TOTAL_PAGES = 17


# ----- PART 1 -----
SEC1 = (1, '时代定调')

def p1_macro():
    slide = new_content_slide(*SEC1, 1, TOTAL_PAGES,
        '宏观背景解读：「十五五」的战略坐标',
        '经济结构转型的^^承上启下^^窗口期 — 规划逻辑与对房地产的传导')
    items = [
        ('高质量发展', '从规模速度型 → 质量效益型；GDP 增速服从于全要素生产率提升'),
        ('新质生产力', '以科技创新为内核；战略性新兴产业 + 未来产业 双轮驱动'),
        ('现代化产业体系', '实体经济 + 数字经济深度融合；安全可控的产业链供应链'),
        ('城乡统筹与^^人口-产业-空间^^再匹配', '都市圈 / 城市群 / 县域经济协同重构空间需求'),
    ]
    top = Inches(2.45)
    for i, (k, v) in enumerate(items):
        y = top + Inches(0.85) * i
        add_rect(slide, Inches(0.78), y, Inches(0.18), Inches(0.7), GOLD)
        add_textbox(slide, Inches(1.05), y + Inches(0.02),
                    Inches(6.5), Inches(0.4),
                    k, size=14, bold=True, color=NAVY)
        add_textbox(slide, Inches(1.05), y + Inches(0.4),
                    Inches(6.5), Inches(0.4),
                    v, size=11, color=TEXT_BODY)

    add_rect(slide, Inches(8.0), Inches(2.3), Inches(4.7), Inches(4.5),
             BG_LIGHT, line=GRAY_LINE)
    add_textbox(slide, Inches(8.2), Inches(2.4), Inches(4.4), Inches(0.4),
                '▎ 政策 → 行业的传导链', size=12, bold=True, color=NAVY)
    chain = [
        ('国家战略', '新质生产力 / 实体经济 / 安全发展', NAVY_LIGHT),
        ('城市发展', '都市圈协同 / 中心城市能级提升', TEAL),
        ('产业空间', '科创园区 / 总部 / 产业社区', ORANGE),
        ('不动产新模式', '存量运营 + 科技服务 + 资产管理', GOLD),
    ]
    cur = Inches(2.95)
    for i, (k, v, c) in enumerate(chain):
        add_rect(slide, Inches(8.25), cur, Inches(4.2), Inches(0.7), WHITE, line=GRAY_LINE)
        add_rect(slide, Inches(8.25), cur, Inches(0.12), Inches(0.7), c)
        add_textbox(slide, Inches(8.5), cur + Inches(0.06),
                    Inches(4.0), Inches(0.3),
                    k, size=11, bold=True, color=c)
        add_textbox(slide, Inches(8.5), cur + Inches(0.36),
                    Inches(4.0), Inches(0.3),
                    v, size=10, color=TEXT_BODY)
        cur += Inches(0.78)
        if i < len(chain) - 1:
            add_textbox(slide, Inches(10.2), cur - Inches(0.13),
                        Inches(0.4), Inches(0.2),
                        '▼', size=10, color=TEXT_MUTE, align=PP_ALIGN.CENTER)


def p1_oldmodel():
    slide = new_content_slide(*SEC1, 2, TOTAL_PAGES,
        '告别旧模式：从^^三高^^驱动到^^稳健提质^^',
        '高杠杆 · 高周转 · 高负债 → 稳健经营 · 提质增效 · 长期主义')

    left_x  = Inches(0.78)
    right_x = Inches(7.0)
    col_w   = Inches(5.7)
    top     = Inches(2.4)

    add_rect(slide, left_x, top, col_w, Inches(0.55), NAVY)
    add_textbox(slide, left_x, top + Inches(0.1), col_w, Inches(0.4),
                '旧模式  ·  规模驱动 (2003-2021)',
                size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_rect(slide, right_x, top, col_w, Inches(0.55), GOLD)
    add_textbox(slide, right_x, top + Inches(0.1), col_w, Inches(0.4),
                '新模式  ·  质量驱动 (十四五末 — 十五五)',
                size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    rows = [
        ('商业逻辑', '土地金融 + 预售制资金循环', '经营性现金流 + 不动产服务'),
        ('增长方式', '扩张土储 / 高周转去化', '存量盘活 / 资产管理 / REITs'),
        ('产品形态', '标准化住宅大宗交付', '好房子 / 好社区 / 好服务'),
        ('收入结构', '开发销售为主', '开发 + 持有 + 运营 + 科技服务'),
        ('风险特征', '高杠杆 · 高负债 · 期限错配', '稳健杠杆 · 长期资金 · 净资产经营'),
        ('竞争核心', '拿地能力 / 融资能力', '运营能力 / 科技能力 / 生态整合'),
    ]
    row_h = Inches(0.55)
    cur = top + Inches(0.55)
    for i, (k, a, b) in enumerate(rows):
        bg = BG_LIGHT if i % 2 == 0 else WHITE
        add_rect(slide, left_x, cur, col_w, row_h, bg, line=GRAY_LINE)
        add_rect(slide, right_x, cur, col_w, row_h, bg, line=GRAY_LINE)
        add_rect(slide, left_x, cur, Inches(1.4), row_h, NAVY)
        add_textbox(slide, left_x, cur + Inches(0.13),
                    Inches(1.4), Inches(0.3),
                    k, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, left_x + Inches(1.55), cur + Inches(0.13),
                    col_w - Inches(1.65), Inches(0.3),
                    a, size=11, color=TEXT_BODY)
        add_textbox(slide, right_x + Inches(0.2), cur + Inches(0.13),
                    col_w - Inches(0.3), Inches(0.3),
                    b, size=11, color=NAVY, bold=True)
        cur += row_h

    add_rect(slide, Inches(0.78), Inches(6.4), Inches(11.77), Inches(0.45),
             NAVY_LIGHT)
    add_textbox(slide, Inches(0.78), Inches(6.46), Inches(11.77), Inches(0.4),
                '▎  本质转换：从^^卖房子^^到^^经营空间、运营人、服务产业^^',
                size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


def p1_thesis():
    slide = new_content_slide(*SEC1, 3, TOTAL_PAGES,
        '核心论点：^^十五五^^竞争维度的根本切换',
        '不再比^^谁拿地多、卖得快^^，而是比^^谁运营得精、谁链接得深^^')

    add_rect(slide, Inches(2.5), Inches(2.3), Inches(8.3), Inches(1.3), NAVY)
    add_rect(slide, Inches(2.5), Inches(2.3), Inches(0.15), Inches(1.3), GOLD)
    add_textbox(slide, Inches(2.8), Inches(2.45), Inches(8.0), Inches(0.5),
                '核心论点 · CORE  THESIS',
                size=11, bold=True, color=GOLD)
    add_textbox(slide, Inches(2.8), Inches(2.8), Inches(8.0), Inches(0.7),
                '^^十五五^^期间的房地产高质量发展 ——',
                size=14, color=WHITE)
    add_textbox(slide, Inches(2.8), Inches(3.15), Inches(8.0), Inches(0.4),
                '= 存量资产的精细化运营  ×  科技产业化的深度融合',
                size=14, bold=True, color=GOLD)

    sub = [
        ('从「土储」到「资产组合」', 'GFA 不再为王；NOI、IRR、Cap Rate 成为通用语言；REITs 打通投融管退闭环。', TEAL),
        ('从「卖空间」到「卖服务」', '招商 / 产业 / 数字 / 金融 服务收入占比上升；空间作为^^流量入口^^。', ORANGE),
        ('从「重资产」到「重平台」', '管理输出 + 科技输出 + 生态输出，平台化能力定义新的护城河。', GOLD),
    ]
    top = Inches(4.05)
    box_w = Inches(3.7)
    gap = Inches(0.18)
    for i, (t, d, c) in enumerate(sub):
        x = Inches(0.78) + (box_w + gap) * i
        add_rect(slide, x, top, box_w, Inches(2.4), WHITE, line=GRAY_LINE)
        add_rect(slide, x, top, box_w, Inches(0.08), c)
        add_textbox(slide, x + Inches(0.3), top + Inches(0.3),
                    Inches(0.6), Inches(0.5),
                    f'0{i+1}', size=22, bold=True, color=c)
        add_textbox(slide, x + Inches(0.3), top + Inches(0.85),
                    box_w - Inches(0.5), Inches(0.5),
                    t, size=14, bold=True, color=NAVY)
        add_rect(slide, x + Inches(0.3), top + Inches(1.35),
                 Inches(0.4), Emu(19050), GOLD)
        add_textbox(slide, x + Inches(0.3), top + Inches(1.5),
                    box_w - Inches(0.5), Inches(0.85),
                    d, size=11, color=TEXT_BODY, line_spacing=1.4)


# ----- PART 2 -----
SEC2 = (2, '破局寻路')

def p2_carrier():
    slide = new_content_slide(*SEC2, 4, TOTAL_PAGES,
        '高质量发展的空间载体演变',
        '从^^居住与办公空间^^  →  ^^科创生态与产业聚集地^^')

    t_top = Inches(2.45)
    add_rect(slide, Inches(0.9), t_top + Inches(0.95),
             Inches(11.5), Emu(28575), GRAY_LINE)
    stages = [
        ('1.0 居住时代', '1998–2010', '解决^^有没有^^\n商品住宅大规模供给', NAVY_LIGHT),
        ('2.0 商办时代', '2010–2020', '解决^^够不够^^\n写字楼 / 商业综合体爆发', TEAL),
        ('3.0 产业时代', '2020–2025', '解决^^专不专^^\n产业园 / 总部 / 产城融合', ORANGE),
        ('4.0 生态时代', '2025+ ^^十五五^^', '解决^^活不活^^\n科创生态 / 平台型不动产', GOLD),
    ]
    n = len(stages)
    seg_w = Inches(11.5) / n
    for i, (name, year, desc, c) in enumerate(stages):
        cx = Inches(0.9) + seg_w * i + seg_w / 2
        add_rect(slide, cx - Inches(0.12),
                 t_top + Inches(0.83), Inches(0.24), Inches(0.24),
                 c, shape=MSO_SHAPE.OVAL)
        add_textbox(slide, cx - Inches(1.1), t_top,
                    Inches(2.2), Inches(0.4),
                    name, size=13, bold=True, color=c, align=PP_ALIGN.CENTER)
        add_textbox(slide, cx - Inches(1.1), t_top + Inches(0.4),
                    Inches(2.2), Inches(0.35),
                    year, size=10, color=TEXT_MUTE, align=PP_ALIGN.CENTER)
        add_textbox(slide, cx - Inches(1.3), t_top + Inches(1.25),
                    Inches(2.6), Inches(1.0),
                    desc, size=11, color=TEXT_BODY,
                    align=PP_ALIGN.CENTER, line_spacing=1.45)

    bottom_top = Inches(4.85)
    add_rect(slide, Inches(0.78), bottom_top, Inches(11.77), Inches(0.45), NAVY)
    add_textbox(slide, Inches(0.78), bottom_top + Inches(0.06),
                Inches(11.77), Inches(0.35),
                '▎ ^^十五五^^时期的不动产载体形态：四种正在崛起的^^高质量空间^^',
                size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    cards = [
        ('科创园区 2.0', '硬科技产线 +\n概念验证中心 + 中试平台'),
        ('产业总部社区', '总部办公 +\n人才公寓 + 配套服务一体化'),
        ('城市更新单元', '老厂房 / 老商业焕新\n承载新消费、新文化、新办公'),
        ('REITs 化资产', '底层资产标准化 +\n公募 REITs 持有运营'),
    ]
    cw = Inches(2.85)
    cy = bottom_top + Inches(0.55)
    for i, (t, d) in enumerate(cards):
        x = Inches(0.78) + (cw + Inches(0.07)) * i
        add_rect(slide, x, cy, cw, Inches(1.4), BG_LIGHT, line=GRAY_LINE)
        add_rect(slide, x, cy, Inches(0.06), Inches(1.4), GOLD)
        add_textbox(slide, x + Inches(0.2), cy + Inches(0.15),
                    cw - Inches(0.3), Inches(0.4),
                    t, size=13, bold=True, color=NAVY)
        add_textbox(slide, x + Inches(0.2), cy + Inches(0.55),
                    cw - Inches(0.3), Inches(0.85),
                    d, size=10.5, color=TEXT_BODY, line_spacing=1.5)


def p2_light_heavy():
    slide = new_content_slide(*SEC2, 5, TOTAL_PAGES,
        '新模式特征 ① 轻重分离：从^^持有规模^^到^^管理半径^^',
        '重资产持有  →  轻资产管理输出 / 品牌输出 / 运营输出')

    left_x = Inches(0.78); right_x = Inches(7.0); col_w = Inches(5.7)
    top = Inches(2.4); h = Inches(4.4)
    add_rect(slide, left_x, top, col_w, h, BG_LIGHT, line=GRAY_LINE)
    add_rect(slide, left_x, top, col_w, Inches(0.55), NAVY)
    add_textbox(slide, left_x, top + Inches(0.1), col_w, Inches(0.4),
                '重 · 资本平台',
                size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    items_h = [
        '底层逻辑：资金 + 资产 + 风控',
        '关键能力：长期资金获取、资产配置、组合管理',
        '代表产品：基础设施 REITs / 持有型物业基金 / 保租房',
        '估值锚点：NOI · Cap Rate · DPU',
        '组织形态：投资委员会 + 资管中后台 + 风控合规',
        '对从业者要求：财务建模、估值、法务结构、跨境资本',
    ]
    cur = top + Inches(0.85)
    for it in items_h:
        add_textbox(slide, left_x + Inches(0.35), cur,
                    col_w - Inches(0.5), Inches(0.45),
                    '●  ' + it, size=11.5, color=TEXT_BODY, line_spacing=1.35)
        cur += Inches(0.55)

    add_rect(slide, right_x, top, col_w, h, BG_LIGHT, line=GRAY_LINE)
    add_rect(slide, right_x, top, col_w, Inches(0.55), GOLD)
    add_textbox(slide, right_x, top + Inches(0.1), col_w, Inches(0.4),
                '轻 · 服务平台',
                size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    items_l = [
        '底层逻辑：品牌 + 标准 + 数据',
        '关键能力：运营 SOP、数字化中台、生态整合',
        '代表产品：商管 / 物管 / 园区运营 / 长租公寓品牌输出',
        '估值锚点：在管面积 · ARPU · 续约率',
        '组织形态：产品中心 + 运营中台 + 客户成功团队',
        '对从业者要求：用户洞察、运营算账、AI 工具、产业资源',
    ]
    cur = top + Inches(0.85)
    for it in items_l:
        add_textbox(slide, right_x + Inches(0.35), cur,
                    col_w - Inches(0.5), Inches(0.45),
                    '●  ' + it, size=11.5, color=NAVY, line_spacing=1.35)
        cur += Inches(0.55)

    add_rect(slide, Inches(5.7), top + h - Inches(0.6),
             Inches(1.95), Inches(0.45), GOLD)
    add_textbox(slide, Inches(5.7), top + h - Inches(0.55),
                Inches(1.95), Inches(0.35),
                '「资本 ⇄ 服务」双轮', size=11, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)


def p2_chancheng():
    slide = new_content_slide(*SEC2, 6, TOTAL_PAGES,
        '新模式特征 ② 产城融合：不动产作为^^超级会客厅^^',
        '聚焦实体经济，链接高校智库 × 科技企业 × 国际创新资源')

    cx, cy = Inches(6.65), Inches(4.55)
    r = Inches(1.0)
    add_rect(slide, cx - r, cy - r, r * 2, r * 2, NAVY,
             shape=MSO_SHAPE.OVAL)
    add_textbox(slide, cx - r, cy - Inches(0.45),
                r * 2, Inches(0.4),
                '不动产', size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(slide, cx - r, cy - Inches(0.05),
                r * 2, Inches(0.4),
                '超级会客厅', size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, cx - r, cy + Inches(0.35),
                r * 2, Inches(0.4),
                'Super  Hub', size=10, color=GOLD, align=PP_ALIGN.CENTER)

    satellites = [
        ('高校 / 智库', '中科院系所、双一流高校实验室、产业研究院',
         Inches(1.2),  Inches(2.95), TEAL),
        ('科技企业', '专精特新、独角兽、上市科技公司、央国企科创平台',
         Inches(9.0),  Inches(2.95), ORANGE),
        ('国际创新资源', '海外加速器、跨境基金、国际行业协会、海归创业团队',
         Inches(1.2),  Inches(5.55), GOLD),
        ('政府 / 园区平台', '区科委、产业引导基金、海关、税务、监管沙盒',
         Inches(9.0),  Inches(5.55), NAVY_LIGHT),
    ]
    for name, desc, x, y, c in satellites:
        w = Inches(3.2); h = Inches(1.4)
        add_rect(slide, x, y, w, h, WHITE, line=GRAY_LINE)
        add_rect(slide, x, y, Inches(0.08), h, c)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.15),
                    w - Inches(0.3), Inches(0.4),
                    name, size=13, bold=True, color=c)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.55),
                    w - Inches(0.3), Inches(0.85),
                    desc, size=10.5, color=TEXT_BODY, line_spacing=1.4)

    add_rect(slide, Inches(0.78), Inches(2.3), Inches(11.77), Inches(0.5),
             BG_LIGHT, line=GRAY_LINE)
    add_textbox(slide, Inches(0.78), Inches(2.37), Inches(11.77), Inches(0.4),
                '▎  空间不再是^^出租物业^^，而是^^价值连接器^^ —— 把人、产业、资本、政策汇聚在同一物理场域',
                size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


def p2_quality_std():
    slide = new_content_slide(*SEC2, 7, TOTAL_PAGES,
        '高质量的^^四把尺子^^：可量化的新模式评价标准',
        '用财务、客户、产业、生态四维度，替代过往单一的销售排名')

    quads = [
        ('F · 财务质量', '经营性现金流 / NOI 增长率\n资产负债率 / 利息保障倍数\nROIC、IRR、Cap Rate', NAVY_LIGHT),
        ('C · 客户质量', '续约率 / NPS / 入驻满意度\n客户结构 (头部 / 腰部 / 长尾)\n生命周期价值 LTV', TEAL),
        ('I · 产业质量', '入驻企业研发强度\n专精特新 / 高新技术企业占比\n产值密度 / 税收密度', ORANGE),
        ('E · 生态质量', '活动 / 路演 / 沙龙频次\n跨界合作机构数量\n投融资促成金额', GOLD),
    ]
    top = Inches(2.45)
    cw = Inches(2.85); ch = Inches(4.1); gap = Inches(0.07)
    for i, (t, d, c) in enumerate(quads):
        x = Inches(0.78) + (cw + gap) * i
        add_rect(slide, x, top, cw, ch, WHITE, line=GRAY_LINE)
        add_rect(slide, x, top, cw, Inches(0.7), c)
        add_textbox(slide, x, top + Inches(0.18), cw, Inches(0.4),
                    t, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, x + Inches(0.3), top + Inches(0.95),
                 Inches(0.4), Emu(19050), GOLD)
        add_textbox(slide, x + Inches(0.3), top + Inches(1.15),
                    cw - Inches(0.6), ch - Inches(1.2),
                    d, size=11.5, color=TEXT_BODY, line_spacing=1.7)

    add_rect(slide, Inches(0.78), Inches(6.78), Inches(11.77), Emu(28575), GOLD)
    add_textbox(slide, Inches(0.78), Inches(6.83), Inches(11.77), Inches(0.3),
                '从^^销售冠军^^到^^运营冠军^^—— ^^十五五^^行业排名的底层逻辑正在重写',
                size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ----- PART 3 -----
SEC3 = (3, '核心动能 · 科技驱动')

def p3_overview():
    slide = new_content_slide(*SEC3, 8, TOTAL_PAGES,
        '科技驱动的三层飞轮：管理 · 空间 · 生态',
        'AI、大模型、数据中台与生态运营，正在重塑不动产价值链')

    layers = [
        ('L1  管理提效', 'AI 与数字化重塑\n投资测算 · 资产管理 · 招商引资',
         '对内 ：让人更少、决策更快、错配更少', NAVY_LIGHT),
        ('L2  空间增值', '科创生态赋能\n沙龙 · 路演 · 加速 · 概念验证',
         '对外 ：让客户更精准、租金溢价、续约更稳', TEAL),
        ('L3  跨界协同', '学术 ⇄ 产业 ⇄ 投资\n超级会客厅 · 平台经济',
         '对生态 ：让不动产成为^^价值链接器^^', GOLD),
    ]
    top = Inches(2.45)
    rh = Inches(1.45); gap = Inches(0.12)
    for i, (lv, what, why, c) in enumerate(layers):
        y = top + (rh + gap) * i
        add_rect(slide, Inches(0.78), y, Inches(11.77), rh, WHITE, line=GRAY_LINE)
        add_rect(slide, Inches(0.78), y, Inches(0.18), rh, c)
        add_textbox(slide, Inches(1.2), y + Inches(0.2),
                    Inches(2.4), Inches(0.5),
                    lv, size=20, bold=True, color=c)
        add_textbox(slide, Inches(1.2), y + Inches(0.85),
                    Inches(2.4), Inches(0.45),
                    '—— 飞轮第 ' + str(i+1) + ' 圈',
                    size=10, color=TEXT_MUTE)
        add_textbox(slide, Inches(3.9), y + Inches(0.18),
                    Inches(4.5), Inches(1.15),
                    what, size=12.5, bold=True, color=NAVY, line_spacing=1.45)
        add_rect(slide, Inches(8.5), y + Inches(0.2),
                 Emu(9525), Inches(1.05), GRAY_LINE)
        add_textbox(slide, Inches(8.7), y + Inches(0.18),
                    Inches(3.7), Inches(1.15),
                    why, size=12, color=TEXT_BODY, line_spacing=1.55)

    add_rect(slide, Inches(0.78), Inches(6.78), Inches(11.77), Emu(28575), GOLD)
    add_textbox(slide, Inches(0.78), Inches(6.83),
                Inches(11.77), Inches(0.3),
                'FLYWHEEL：管理提效  →  空间增值  →  跨界协同  →  反哺管理 (数据 + 案例 + 模型迭代)',
                size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


def p3_ai_invest():
    slide = new_content_slide(*SEC3, 9, TOTAL_PAGES,
        '管理提效 ① 大模型 × 投资测算：从^^经验拍脑袋^^到^^模型推演^^',
        '把投决会前的 80% 体力活，交给 LLM + 量化模型')

    left_x = Inches(0.78); col_w = Inches(5.5); top = Inches(2.4)
    add_rect(slide, left_x, top, col_w, Inches(0.55), NAVY)
    add_textbox(slide, left_x, top + Inches(0.1), col_w, Inches(0.4),
                '应用场景  ·  Where to Use',
                size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    scenes = [
        ('城市与板块研判',
         'LLM 抓取规划、人口、产业、交通、消费数据 ;  自动生成板块画像与排序'),
        ('地块 / 资产尽调',
         'Agent 拉取规划条件 / 周边可比 / 限价信息，10 分钟出^^一页纸^^备忘录'),
        ('现金流测算',
         '模型按情景 (中性 / 乐观 / 压力) 自动跑动态 IRR / NPV / 敏感性矩阵'),
        ('投决会问询应答',
         'RAG 检索历史项目 + 行业研报，模拟^^老板会问的 30 个问题^^自检'),
    ]
    cur = top + Inches(0.7)
    for k, v in scenes:
        add_rect(slide, left_x, cur, col_w, Inches(0.95), BG_LIGHT, line=GRAY_LINE)
        add_rect(slide, left_x, cur, Inches(0.08), Inches(0.95), GOLD)
        add_textbox(slide, left_x + Inches(0.25), cur + Inches(0.1),
                    col_w - Inches(0.4), Inches(0.35),
                    k, size=12.5, bold=True, color=NAVY)
        add_textbox(slide, left_x + Inches(0.25), cur + Inches(0.45),
                    col_w - Inches(0.4), Inches(0.5),
                    v, size=10.5, color=TEXT_BODY, line_spacing=1.4)
        cur += Inches(1.05)

    right_x = Inches(6.55); col_w_r = Inches(6.0)
    add_rect(slide, right_x, top, col_w_r, Inches(0.55), GOLD)
    add_textbox(slide, right_x, top + Inches(0.1), col_w_r, Inches(0.4),
                '可量化的价值  ·  Why It Matters',
                size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    metrics = [
        ('70%', '尽调资料整理时间下降'),
        ('10×', '可比案例检索效率提升'),
        ('3 天 → 3 小时', '项目^^一页纸^^备忘录出稿'),
        ('≥5 套', '情景测算并行可视化'),
    ]
    mt = top + Inches(0.7)
    mw = (col_w_r - Inches(0.15)) / 2
    mh = Inches(1.45)
    for i, (n, d) in enumerate(metrics):
        rx = right_x + (mw + Inches(0.15)) * (i % 2)
        ry = mt + (mh + Inches(0.15)) * (i // 2)
        add_rect(slide, rx, ry, mw, mh, WHITE, line=GRAY_LINE)
        add_rect(slide, rx, ry, Inches(0.08), mh, NAVY)
        add_textbox(slide, rx + Inches(0.2), ry + Inches(0.2),
                    mw - Inches(0.3), Inches(0.6),
                    n, size=24, bold=True, color=NAVY)
        add_textbox(slide, rx + Inches(0.2), ry + Inches(0.85),
                    mw - Inches(0.3), Inches(0.5),
                    d, size=11, color=TEXT_BODY, line_spacing=1.4)

    add_rect(slide, right_x, mt + (mh + Inches(0.15)) * 2,
             col_w_r, Inches(0.55), NAVY_LIGHT)
    add_textbox(slide, right_x, mt + (mh + Inches(0.15)) * 2 + Inches(0.08),
                col_w_r, Inches(0.4),
                '▎ 底线：模型只做^^草稿^^，决策仍由人 + 投委会负责；数据合规 / 隐私不可让位',
                size=10.5, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


def p3_ai_assetmgmt():
    slide = new_content_slide(*SEC3, 10, TOTAL_PAGES,
        '管理提效 ② AI × 资产管理：让^^楼^^自己说话',
        '数字化中台 + 大模型 Agent，把^^资产管家^^从被动维护推向主动经营')

    steps = [
        ('感知', 'IoT + BIM\n空调 / 电梯 / 能耗 / 客流 / 摄像头', NAVY_LIGHT),
        ('分析', '数据中台 + 大模型\n异常识别 / 趋势预测 / 客户画像', TEAL),
        ('决策', 'AI Copilot\n租赁定价 / 招商优先级 / 维保排程', ORANGE),
        ('行动', '自动化执行\n工单派发 / 报表生成 / 客户触达', GOLD),
    ]
    top = Inches(2.5)
    sw = Inches(2.85); sh_h = Inches(2.0); gap = Inches(0.15)
    for i, (k, d, c) in enumerate(steps):
        x = Inches(0.78) + (sw + gap) * i
        add_rect(slide, x, top, sw, sh_h, WHITE, line=GRAY_LINE)
        add_rect(slide, x, top, sw, Inches(0.5), c)
        add_textbox(slide, x, top + Inches(0.1), sw, Inches(0.35),
                    k, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), top + Inches(0.7),
                    sw - Inches(0.4), sh_h - Inches(0.8),
                    d, size=11.5, color=TEXT_BODY,
                    align=PP_ALIGN.CENTER, line_spacing=1.5)
        if i < 3:
            add_textbox(slide, x + sw - Inches(0.05), top + Inches(0.85),
                        Inches(0.25), Inches(0.4),
                        '▶', size=14, color=GOLD, align=PP_ALIGN.CENTER)

    bot_top = Inches(4.8)
    add_rect(slide, Inches(0.78), bot_top, Inches(11.77), Inches(0.42), NAVY)
    add_textbox(slide, Inches(0.78), bot_top + Inches(0.05),
                Inches(11.77), Inches(0.32),
                '▎  典型 KPI 改善  ·  从^^管事^^到^^管效益^^',
                size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    kpis = [
        ('能耗', '↓ 8-15%', 'AI 调优 + 错峰策略'),
        ('空置', '↓ 1-3 个月', '客户画像 + 主动招商'),
        ('租金溢价', '↑ 3-8%', '定价模型 + 服务捆绑'),
        ('客户满意度', '↑ 10+ pt', '工单响应 + 服务可视化'),
    ]
    cy = bot_top + Inches(0.55)
    cw = (Inches(11.77) - Inches(0.45)) / 4
    for i, (k, n, d) in enumerate(kpis):
        x = Inches(0.78) + (cw + Inches(0.15)) * i
        add_rect(slide, x, cy, cw, Inches(1.4), BG_LIGHT, line=GRAY_LINE)
        add_textbox(slide, x + Inches(0.2), cy + Inches(0.12),
                    cw - Inches(0.3), Inches(0.35),
                    k, size=12, bold=True, color=NAVY)
        add_textbox(slide, x + Inches(0.2), cy + Inches(0.45),
                    cw - Inches(0.3), Inches(0.5),
                    n, size=20, bold=True, color=ORANGE)
        add_textbox(slide, x + Inches(0.2), cy + Inches(0.95),
                    cw - Inches(0.3), Inches(0.4),
                    d, size=10.5, color=TEXT_BODY, line_spacing=1.35)


def p3_ai_zhaoshang():
    slide = new_content_slide(*SEC3, 11, TOTAL_PAGES,
        '管理提效 ③ 自动化招商引资：从^^扫楼陌拜^^到^^AI 雷达^^',
        '把企业搜寻、画像、触达、跟进的全链路装进数字化中台')

    funnel_levels = [
        ('企业池',         '工商 / 税务 / 招标 / 招聘 / 专利 / 融资数据汇聚', NAVY_LIGHT, 0.95),
        ('画像与匹配',     'LLM 抽取行业、阶段、空间偏好 → 与园区画像匹配', TEAL,        0.78),
        ('线索分级',       'AI 评分 (热 / 温 / 冷) + 接触建议',                ORANGE,       0.62),
        ('自动化触达',     '一键生成 BD 话术 / 邮件 / 提案 / 微信内容',          GOLD,         0.46),
        ('成交与履约',     '电子签 / 入驻系统 / 服务工单 / 续约预测',             NAVY,         0.32),
    ]
    top = Inches(2.5)
    fh = Inches(0.7); gap = Inches(0.1)
    base_w = Inches(7.5)
    for i, (k, d, c, ratio) in enumerate(funnel_levels):
        w = base_w * ratio
        x = Inches(0.9) + (base_w - w) / 2
        y = top + (fh + gap) * i
        add_rect(slide, x, y, w, fh, c)
        add_textbox(slide, x, y + Inches(0.18), w, Inches(0.4),
                    k, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(8.55), y + Inches(0.18),
                    Inches(4.4), Inches(0.4),
                    d, size=11, color=TEXT_BODY)

    bot_top = Inches(6.4)
    add_rect(slide, Inches(0.78), bot_top, Inches(11.77), Inches(0.5),
             BG_LIGHT, line=GRAY_LINE)
    add_textbox(slide, Inches(0.95), bot_top + Inches(0.1),
                Inches(11.5), Inches(0.35),
                '▎ 招商成本/家  ↓ 30-50%   ·   线索→成交 转化  ↑ 2-3×   ·   ^^专精特新^^目标客户命中率显著提升',
                size=11.5, bold=True, color=NAVY)


def p3_space_ecosystem():
    slide = new_content_slide(*SEC3, 12, TOTAL_PAGES,
        '空间增值 · 科创生态赋能：把^^楼宇^^运营成^^创新场^^',
        '通过持续的内容运营，使物理空间成为产业要素的^^高频聚合点^^')

    left_x = Inches(0.78); col_w = Inches(5.6); top = Inches(2.4)
    add_rect(slide, left_x, top, col_w, Inches(0.5), NAVY)
    add_textbox(slide, left_x, top + Inches(0.07), col_w, Inches(0.4),
                '▎  科创生态运营 · ^^四件套^^',
                size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    items = [
        ('前沿科技沙龙', '高校老师 + 一线工程师小范围深聊 ;\n聚焦 AI / 生物医药 / 新能源 / 半导体'),
        ('产业 / 投资路演', '项目方 × 投资人 × 大企业 CVC 三方撮合 ;\n月度路演 + 季度专场'),
        ('国际创新对接', '联合海外加速器 / 协会 ;\nLanding Pad / 软着陆服务 / 跨境服务包'),
        ('企业家私享会', 'Family Office × 上市公司创始人小范围闭门 ;\n输出^^被需要^^的稀缺连接'),
    ]
    cur = top + Inches(0.65)
    for k, v in items:
        add_rect(slide, left_x, cur, col_w, Inches(0.95), WHITE, line=GRAY_LINE)
        add_rect(slide, left_x, cur, Inches(0.06), Inches(0.95), GOLD)
        add_textbox(slide, left_x + Inches(0.2), cur + Inches(0.1),
                    col_w - Inches(0.3), Inches(0.3),
                    k, size=12.5, bold=True, color=NAVY)
        add_textbox(slide, left_x + Inches(0.2), cur + Inches(0.4),
                    col_w - Inches(0.3), Inches(0.5),
                    v, size=10.5, color=TEXT_BODY, line_spacing=1.4)
        cur += Inches(1.05)

    right_x = Inches(6.65); col_w_r = Inches(6.0)
    add_rect(slide, right_x, top, col_w_r, Inches(0.5), GOLD)
    add_textbox(slide, right_x, top + Inches(0.07), col_w_r, Inches(0.4),
                '▎  从^^内容^^到^^租金^^的价值传导',
                size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    chain = [
        ('高频内容',        '一年 ≥ 50 场高质量活动',            NAVY_LIGHT),
        ('精准客流',        '目标客户主动到访 / 复访',           TEAL),
        ('品牌势能',        '成为细分领域的^^必到地标^^',         ORANGE),
        ('溢价与黏性',      '租金溢价 + 续约率 + 转介绍提升',   GOLD),
        ('资产估值',        'NOI 提升 → Cap Rate 收敛 → 估值上行', NAVY),
    ]
    cur = top + Inches(0.65)
    for i, (k, v, c) in enumerate(chain):
        add_rect(slide, right_x, cur, col_w_r, Inches(0.7),
                 WHITE, line=GRAY_LINE)
        add_rect(slide, right_x, cur, Inches(0.1), Inches(0.7), c)
        add_textbox(slide, right_x + Inches(0.25), cur + Inches(0.18),
                    Inches(2.0), Inches(0.4),
                    k, size=12, bold=True, color=c)
        add_textbox(slide, right_x + Inches(2.4), cur + Inches(0.22),
                    col_w_r - Inches(2.5), Inches(0.4),
                    v, size=11, color=TEXT_BODY)
        cur += Inches(0.78)


def p3_zjtx_opc():
    slide = new_content_slide(*SEC3, 13, TOTAL_PAGES,
        '精准服务两类新兴主体：「专精特新」与「一人公司 OPC」',
        '供给侧匹配新主体——空间、服务、资本、合规一体化')

    top1 = Inches(2.4)
    add_rect(slide, Inches(0.78), top1, Inches(11.77), Inches(2.0),
             WHITE, line=GRAY_LINE)
    add_rect(slide, Inches(0.78), top1, Inches(0.18), Inches(2.0), TEAL)
    add_textbox(slide, Inches(1.05), top1 + Inches(0.15),
                Inches(11), Inches(0.4),
                '对象 ① ：「专精特新」中小企业  ·  Specialized · Sophisticated · Featured · Novel',
                size=14, bold=True, color=NAVY)
    needs1 = [
        ('空间需求', '中小型独立单元 + 共享中试 / 实验空间 + 弹性扩容'),
        ('服务需求', '知识产权、检验检测、人才公寓、子女教育、专项补贴申报'),
        ('资本需求', '种子 / 天使 / Pre-A 衔接 ; 与本地引导基金 / CVC 撮合'),
        ('数据需求', '行业景气、政策变动、招投标信息的定向推送'),
    ]
    for i, (k, v) in enumerate(needs1):
        x = Inches(1.05) + (Inches(2.85) + Inches(0.05)) * i
        add_rect(slide, x, top1 + Inches(0.7), Inches(2.85), Inches(1.2),
                 BG_LIGHT, line=GRAY_LINE)
        add_textbox(slide, x + Inches(0.15), top1 + Inches(0.78),
                    Inches(2.6), Inches(0.35),
                    k, size=12, bold=True, color=TEAL)
        add_textbox(slide, x + Inches(0.15), top1 + Inches(1.12),
                    Inches(2.6), Inches(0.75),
                    v, size=10.5, color=TEXT_BODY, line_spacing=1.45)

    top2 = Inches(4.55)
    add_rect(slide, Inches(0.78), top2, Inches(11.77), Inches(2.45),
             WHITE, line=GRAY_LINE)
    add_rect(slide, Inches(0.78), top2, Inches(0.18), Inches(2.45), ORANGE)
    add_textbox(slide, Inches(1.05), top2 + Inches(0.15),
                Inches(11.5), Inches(0.4),
                '对象 ② ：「一人公司 / OPC · One Person Company」与超级个体户  ·  AI 时代的新经济细胞',
                size=14, bold=True, color=NAVY)
    add_textbox(slide, Inches(1.05), top2 + Inches(0.55),
                Inches(11.5), Inches(0.4),
                '代表：独立顾问、AI 工程师、内容创作者、自由律师、跨境电商主理人、独立基金 GP',
                size=11, color=TEXT_MUTE)
    needs2 = [
        ('产品形态', '工位 + 私密会议 + 直播间 + 内容工坊\n按月 / 按日 / 按小时灵活计费'),
        ('配套服务', '工商代办 / 财税 / 法律 / 出海 / 支付通道\n^^开公司像开会员^^一样轻'),
        ('社群运营', '小型 Mastermind 圈层 / 行业小组\n让独行者也有^^组织^^'),
        ('空间美学', '可拍可发可分享的^^内容友好型^^空间\n品牌即流量入口'),
    ]
    for i, (k, v) in enumerate(needs2):
        x = Inches(1.05) + (Inches(2.85) + Inches(0.05)) * i
        add_rect(slide, x, top2 + Inches(1.0), Inches(2.85), Inches(1.35),
                 BG_LIGHT, line=GRAY_LINE)
        add_textbox(slide, x + Inches(0.15), top2 + Inches(1.08),
                    Inches(2.6), Inches(0.35),
                    k, size=12, bold=True, color=ORANGE)
        add_textbox(slide, x + Inches(0.15), top2 + Inches(1.42),
                    Inches(2.6), Inches(0.9),
                    v, size=10.5, color=TEXT_BODY, line_spacing=1.45)


def p3_crossover():
    slide = new_content_slide(*SEC3, 14, TOTAL_PAGES,
        '跨界协同：打通「学术圈 ⇄ 产业圈 ⇄ 投资圈」',
        '三圈打通，不动产从^^租赁中介^^升维为^^产业服务平台^^')

    cy = Inches(4.4)
    r = Inches(1.3)
    centers = [
        (Inches(5.3),  cy, '学术圈', '高校 · 院所 · 智库',  TEAL),
        (Inches(9.7),  cy, '产业圈', '企业 · 链主 · 链群',  ORANGE),
        (Inches(7.5),  cy + Inches(1.3), '投资圈', 'VC / PE · CVC · 引导基金', GOLD),
    ]
    for cx, cy_, name, sub, c in centers:
        add_rect(slide, cx - r, cy_ - r, r * 2, r * 2, c,
                 shape=MSO_SHAPE.OVAL)
        add_textbox(slide, cx - r, cy_ - Inches(0.25),
                    r * 2, Inches(0.4),
                    name, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, cx - r, cy_ + Inches(0.18),
                    r * 2, Inches(0.4),
                    sub, size=10, color=WHITE, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(6.8), cy + Inches(0.25),
             Inches(1.4), Inches(0.7), NAVY)
    add_textbox(slide, Inches(6.8), cy + Inches(0.4),
                Inches(1.4), Inches(0.4),
                '不动产平台', size=11, bold=True, color=GOLD,
                align=PP_ALIGN.CENTER)

    left_x = Inches(0.78); col_w = Inches(3.8); top = Inches(2.4)
    add_rect(slide, left_x, top, col_w, Inches(4.7),
             BG_LIGHT, line=GRAY_LINE)
    add_textbox(slide, left_x + Inches(0.2), top + Inches(0.2),
                col_w - Inches(0.3), Inches(0.4),
                '▎ 协同机制', size=13, bold=True, color=NAVY)
    mechs = [
        ('成果转化通道', '高校实验室 → 概念验证中心 → 中试基地 → 产线'),
        ('产学研联合体', '企业出题、高校解题、平台搭台、基金接力'),
        ('人才环流', '教授驻企 / 工程师入校 / 学生入企实习'),
        ('数据 & 案例库', '把^^链接^^沉淀成可复用的资产 (而非个人人脉)'),
    ]
    cur = top + Inches(0.7)
    for k, v in mechs:
        add_rect(slide, left_x + Inches(0.15), cur,
                 Inches(0.06), Inches(0.95), GOLD)
        add_textbox(slide, left_x + Inches(0.3), cur,
                    col_w - Inches(0.4), Inches(0.3),
                    k, size=11.5, bold=True, color=NAVY)
        add_textbox(slide, left_x + Inches(0.3), cur + Inches(0.32),
                    col_w - Inches(0.4), Inches(0.7),
                    v, size=10.5, color=TEXT_BODY, line_spacing=1.4)
        cur += Inches(1.0)

    add_rect(slide, Inches(0.78), Inches(6.65), Inches(11.77), Inches(0.4), NAVY)
    add_textbox(slide, Inches(0.78), Inches(6.71), Inches(11.77), Inches(0.35),
                '▎ 不动产价值的^^二次跃升^^ = 物理空间的租金价值 + 链接资源的服务价值',
                size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


def p3_case_matrix():
    slide = new_content_slide(*SEC3, 15, TOTAL_PAGES,
        '实战案例矩阵：科技 × 不动产的六种可复制范式',
        '每个范式都对应^^问题—工具—成效^^三段式落地')

    cases = [
        ('AI 投决助手',     '投资测算 / 板块研判',  '70% 案头时间释放，3 天 → 3 小时出备忘录', NAVY_LIGHT),
        ('智能资管中台',   '持有型物业经营',        '能耗 ↓10%，空置月数 ↓30%', TEAL),
        ('AI 招商雷达',     '园区 / 写字楼招商',    '线索成本 ↓40%，转化率 2-3×', ORANGE),
        ('科创内容工厂',   '园区生态运营',          '年 ≥50 场，租金溢价 5-8%', GOLD),
        ('OPC 服务包',      '灵活办公 + 一人公司',   '客户结构多元、续费率显著提升', NAVY),
        ('REITs 资管闭环',  '持有型资产证券化',      '估值锚点切换为 NOI / DPU', NAVY_LIGHT),
    ]
    top = Inches(2.4)
    cw = Inches(3.85); ch = Inches(2.05); gap_x = Inches(0.1); gap_y = Inches(0.15)
    for i, (name, scene, effect, c) in enumerate(cases):
        col = i % 3; row = i // 3
        x = Inches(0.78) + (cw + gap_x) * col
        y = top + (ch + gap_y) * row
        add_rect(slide, x, y, cw, ch, WHITE, line=GRAY_LINE)
        add_rect(slide, x, y, cw, Inches(0.45), c)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.08),
                    cw - Inches(0.3), Inches(0.35),
                    f'范式 {i+1:02d}  ·  {name}',
                    size=12, bold=True, color=WHITE)
        add_textbox(slide, x + Inches(0.25), y + Inches(0.6),
                    cw - Inches(0.4), Inches(0.4),
                    '应用场景', size=10, color=TEXT_MUTE, bold=True)
        add_textbox(slide, x + Inches(0.25), y + Inches(0.92),
                    cw - Inches(0.4), Inches(0.4),
                    scene, size=12, color=NAVY, bold=True)
        add_rect(slide, x + Inches(0.25), y + Inches(1.32),
                 cw - Inches(0.5), Emu(9525), GRAY_LINE)
        add_textbox(slide, x + Inches(0.25), y + Inches(1.4),
                    cw - Inches(0.4), Inches(0.3),
                    '可量化成效', size=10, color=TEXT_MUTE, bold=True)
        add_textbox(slide, x + Inches(0.25), y + Inches(1.62),
                    cw - Inches(0.4), Inches(0.4),
                    effect, size=10.5, color=TEXT_BODY, line_spacing=1.4)


# ----- PART 4 -----
SEC4 = (4, '总结展望')

def p4_consensus():
    slide = new_content_slide(*SEC4, 16, TOTAL_PAGES,
        '发展共识：理论与实践，缺一不可',
        '学术界提供^^为什么^^，实业界回答^^怎么做^^')

    cols = [
        ('学 · 学术界',
         'ACADEMIA  ·  Why & What',
         [
             '宏观研判 ：经济周期、人口结构、产业升级方向',
             '理论框架 ：新模式 / 新质生产力 / 新城市化',
             '数据洞察 ：长周期数据、跨国比较、政策评估',
             '人才培养 ：财务 + 产业 + 数字复合人才',
         ],
         NAVY_LIGHT, TEAL),
        ('产 · 实业界',
         'INDUSTRY  ·  How & Now',
         [
             '蹚路探索 ：在真实项目里验证模型与产品',
             '工具落地 ：用 AI / 数字化把方法变成生产力',
             '客户回馈 ：以一线痛点反向定义研究议题',
             '生态搭建 ：以平台撬动学术与资本资源',
         ],
         GOLD, ORANGE),
    ]
    top = Inches(2.4); ch = Inches(4.2)
    for i, (k, en, lines, c1, c2) in enumerate(cols):
        x = Inches(0.78) + (Inches(5.95) + Inches(0.15)) * i
        w = Inches(5.95)
        add_rect(slide, x, top, w, ch, WHITE, line=GRAY_LINE)
        add_rect(slide, x, top, w, Inches(0.7), c1)
        add_textbox(slide, x + Inches(0.3), top + Inches(0.12),
                    w - Inches(0.5), Inches(0.4),
                    k, size=18, bold=True, color=WHITE)
        add_textbox(slide, x + Inches(0.3), top + Inches(0.42),
                    w - Inches(0.5), Inches(0.3),
                    en, size=10, color=GOLD, bold=True)
        cur = top + Inches(0.95)
        for ln in lines:
            add_rect(slide, x + Inches(0.3), cur + Inches(0.1),
                     Inches(0.08), Inches(0.7), c2)
            add_textbox(slide, x + Inches(0.5), cur,
                        w - Inches(0.7), Inches(0.85),
                        ln, size=12, color=TEXT_BODY, line_spacing=1.45)
            cur += Inches(0.85)

    add_rect(slide, Inches(0.78), Inches(6.78), Inches(11.77), Emu(28575), GOLD)
    add_textbox(slide, Inches(0.78), Inches(6.83), Inches(11.77), Inches(0.3),
                '学术指引方向，实业蹚出新路 —— 这正是本次跨界汇报的意义所在',
                size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


def p4_action():
    slide = new_content_slide(*SEC4, 17, TOTAL_PAGES,
        '^^十五五^^从业者行动清单：知 与 行',
        '对从业者的 6 条建议 —— ^^今晚就能开始做的事^^')

    actions = [
        ('【知】重读规划',     '^^十五五^^规划纲要 + 行业专项规划 + 城市层面^^四个一^^文件',         NAVY_LIGHT, '认知'),
        ('【知】重学财务',     '把 NOI / Cap Rate / IRR / DPU 变成日常通用语言',                    TEAL,        '认知'),
        ('【知】重看产业',     '选定 1-2 个产业 (如 AI / 半导体 / 创新药) 持续跟踪，建立深度',      ORANGE,      '认知'),
        ('【行】用一次大模型', '本周用 LLM 做一份^^项目一页纸^^或^^客户画像^^，亲身体感',           GOLD,        '行动'),
        ('【行】跑一次活动',   '组织一场^^小而精^^的科创沙龙 / 路演，验证生态运营能力',              NAVY,        '行动'),
        ('【行】算一次资产账', '选一栋楼，按^^NOI 经营^^重新算一遍 5 年现金流，跳出销售思维',        NAVY_LIGHT,  '行动'),
    ]
    top = Inches(2.35)
    rh = Inches(0.62); gap = Inches(0.1)
    for i, (k, v, c, tag) in enumerate(actions):
        y = top + (rh + gap) * i
        add_rect(slide, Inches(0.78), y, Inches(11.77), rh,
                 WHITE, line=GRAY_LINE)
        add_rect(slide, Inches(0.78), y, Inches(0.12), rh, c)
        add_textbox(slide, Inches(0.95), y + Inches(0.15),
                    Inches(2.4), Inches(0.4),
                    k, size=13, bold=True, color=NAVY)
        add_rect(slide, Inches(3.45), y + Inches(0.13),
                 Emu(9525), Inches(0.36), GRAY_LINE)
        add_textbox(slide, Inches(3.6), y + Inches(0.15),
                    Inches(7.5), Inches(0.4),
                    v, size=11.5, color=TEXT_BODY)
        add_rect(slide, Inches(11.2), y + Inches(0.13),
                 Inches(1.2), Inches(0.36), c)
        add_textbox(slide, Inches(11.2), y + Inches(0.17),
                    Inches(1.2), Inches(0.35),
                    tag, size=10, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)


def p4_closing():
    slide = slide_blank(prs, bg=NAVY)
    add_rect(slide, 0, 0, Inches(0.4), SH, GOLD)
    add_textbox(slide, Inches(0.9), Inches(0.7), Inches(8), Inches(0.4),
                'CLOSING  ·  和 光 同 尘',
                size=12, bold=True, color=GOLD)

    add_textbox(slide, Inches(0.9), Inches(1.4), Inches(11.5), Inches(1.2),
                '结语：行业在阵痛，但土壤在生长',
                size=34, bold=True, color=WHITE)

    add_rect(slide, Inches(0.9), Inches(2.65), Inches(0.6), Emu(28575), GOLD)
    add_textbox(slide, Inches(0.9), Inches(2.8), Inches(11.5), Inches(0.55),
                '和光同尘  ·  顺应科技大势与国家宏观规划',
                size=18, color=GOLD)

    paragraphs = [
        '▍ 房地产正在告别^^高歌猛进^^，但并未告别^^价值创造^^。',
        '▍ ^^十五五^^开启的，是从^^开发驱动^^到^^运营驱动^^、从^^空间红利^^到^^科技红利^^的换挡。',
        '▍ 当不动产真正成为科技、产业、资本、人才的^^高频聚合点^^，',
        '    它就是孕育「新质生产力」最肥沃的土壤之一。',
    ]
    cur = Inches(3.9)
    for p in paragraphs:
        add_textbox(slide, Inches(0.9), cur, Inches(11.5), Inches(0.55),
                    p, size=14, color=WHITE, line_spacing=1.5)
        cur += Inches(0.55)

    add_rect(slide, Inches(0.9), Inches(6.4), Inches(3.0), Inches(0.7), GOLD)
    add_textbox(slide, Inches(0.9), Inches(6.55), Inches(3.0), Inches(0.5),
                'Q  &  A   敬请指正',
                size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(4.1), Inches(6.6), Inches(8.5), Inches(0.5),
                '—  汇报完毕，感谢聆听  —',
                size=14, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)


def build_all():
    build_cover()
    build_toc()

    build_section_divider(
        1, '时代定调：「十五五」开局与房地产周期的新旧转换',
        'PART 01  ·  Era  Setting',
        '宏观政策切入  ·  告别旧模式  ·  提出核心论点',
        [
            ('01', '^^十五五^^战略定位\n经济结构转型坐标'),
            ('02', '告别^^三高^^旧模式\n稳健经营 提质增效'),
            ('03', '高质量发展核心论点\n精细化运营 × 科技融合'),
        ])
    p1_macro()
    p1_oldmodel()
    p1_thesis()

    build_section_divider(
        2, '破局寻路：房地产「新模式」的核心内涵与高质量标准',
        'PART 02  ·  Break  Through',
        '空间载体演变  ·  轻重分离  ·  产城融合  ·  四把尺子',
        [
            ('01', '空间载体进化\n从居住到生态'),
            ('02', '轻重分离 + 产城融合\n双特征解构'),
            ('03', '高质量评价四维\n财务 / 客户 / 产业 / 生态'),
        ])
    p2_carrier()
    p2_light_heavy()
    p2_chancheng()
    p2_quality_std()

    build_section_divider(
        3, '核心动能：科技创新如何驱动不动产高质量发展',
        'PART 03  ·  Tech-Powered  Engine',
        'AI 与大模型  ·  科创生态赋能  ·  跨界协同  ·  实战范式',
        [
            ('01', '管理提效\nAI × 投资 / 资管 / 招商'),
            ('02', '空间增值\n科创生态运营 + 新主体服务'),
            ('03', '跨界协同\n学术 ⇄ 产业 ⇄ 投资'),
        ])
    p3_overview()
    p3_ai_invest()
    p3_ai_assetmgmt()
    p3_ai_zhaoshang()
    p3_space_ecosystem()
    p3_zjtx_opc()
    p3_crossover()
    p3_case_matrix()

    build_section_divider(
        4, '总结展望：「十五五」时期从业者的知与行',
        'PART 04  ·  Outlook  &  Action',
        '学术与实业的共识  ·  行动清单  ·  和光同尘',
        [
            ('01', '理论与实践共识\n学指方向 实蹚新路'),
            ('02', '六条行动清单\n今晚就能开始的事'),
            ('03', '结语：阵痛中的土壤\n孕育新质生产力'),
        ])
    p4_consensus()
    p4_action()
    p4_closing()


build_all()

out = '/workspace/十五五不动产高质量发展汇报.pptx'
prs.save(out)
print('saved:', out, 'slides:', len(prs.slides))
