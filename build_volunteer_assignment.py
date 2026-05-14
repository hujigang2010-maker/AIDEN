"""
2026 人工智能商业化落地与硬核投资破局峰会
志愿者职能分工 - 生成 Excel 分工表 + PPT
作者：总策划（按本议程定制）
"""

from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy

# =============================================================================
# 通用数据 - 议程
# =============================================================================
EVENT_TITLE = "2026 人工智能商业化落地与硬核投资破局峰会"
EVENT_SUB = "寻找 AI 时代的超级个体、新质资产与资本新风口"
EVENT_DATE = "2026 年 5 月 22 日 13:00 - 20:30"
EVENT_VENUE = "上海 · 北外滩 · 一滴水"

AGENDA = [
    ("13:00-13:30", "嘉宾签到与入场",              "C/B/D",  "签到处、引导动线、媒体合影区"),
    ("13:30-13:35", "开幕致辞（姚志勇）",         "C",      "舞台、主持、上下场"),
    ("13:35-14:00", "主旨演讲1·白硕",              "C/E",    "PPT、麦克风、计时"),
    ("14:00-14:25", "主旨演讲2·夏春",              "C/E",    "PPT、麦克风、计时"),
    ("14:25-15:05", "圆桌一·AI硬核（王珏主持4嘉宾）","C/E",   "圆桌座次、麦克风、计时、互动"),
    ("15:05-15:45", "巅峰对话·云生态（胡继刚主持3嘉宾）","C/E","圆桌座次、麦克风、计时"),
    ("15:45-16:35", "陆家嘴交响乐团·时光音乐会",    "C/F/D",  "乐团对接、灯光、社交茶歇"),
    ("16:35-16:55", "主旨演讲3·王维军",            "C/E",    "PPT、麦克风、计时"),
    ("16:55-17:15", "主旨演讲4·寇文红",            "C/E",    "PPT、麦克风、计时"),
    ("17:15-17:35", "主旨演讲5·张露瑶",            "C/E",    "PPT、麦克风、计时"),
    ("17:35-18:15", "圆桌二·投资（黄欣主持4嘉宾）", "C/E",   "圆桌座次、麦克风、计时"),
    ("18:15-18:30", "颁奖典礼暨闭幕",              "C/F/D",  "奖杯、颁奖嘉宾引导、合影"),
    ("18:30-20:30", "VIP高端闭门晚宴",             "B/F/G",  "席位、餐饮、对接、安保"),
]

# =============================================================================
# 组别架构
# =============================================================================
GROUPS = [
    {
        "code": "A",
        "name": "统筹协调组",
        "leader": "韦佳玉（执行总监）",
        "size": "3 人",
        "color": "1F4E79",
        "duty_summary": "向总策划负责，统筹全场，协同各组进度、对外口径、突发决策",
        "sub_duties": [
            ("总指挥（胡继刚）",       "全场最高决策；与主办/协办方对接；流程总控"),
            ("执行总监+文件总管（韦佳玉）","对接 B/C/D/E/F/G 六组；议程脚本/会议记录/复盘主持；对讲机调度中枢"),
            ("科创资源对接专员（春子）","聚焦科创赞助谈判+科技展位推进（自己提的科技展示建议自己主推）"),
        ],
    },
    {
        "code": "B",
        "name": "嘉宾接待组",
        "leader": "王胜（组长）",
        "size": "5 人",
        "color": "2E75B6",
        "duty_summary": "覆盖嘉宾从邀约确认 → 抵达 → 接送 → 签到 → 引导 → 离场 → 晚宴全链路",
        "sub_duties": [
            ("VIP一对一",       "为主旨演讲嘉宾/主持/特邀机构嘉宾配备 1V1 陪同（白硕、夏春、寇文红等）"),
            ("机场/酒店接送",   "外地嘉宾接机/接站、酒店入住对接、当日接送车辆调度"),
            ("签到台",          "按嘉宾分级（VVIP / 演讲嘉宾 / 机构嘉宾 / 普通嘉宾）分通道签到"),
            ("引导动线",        "电梯口→签到台→主会场→洗手间→茶歇区→晚宴厅 全程引导"),
            ("休息室管理",      "嘉宾休息室饮水、补妆、PPT预演、专属服务"),
            ("名片/合影",       "重要嘉宾合影、名片收集、对接需求登记"),
        ],
    },
    {
        "code": "C",
        "name": "会务执行组",
        "leader": "陈潇 Kelly（组长）",
        "size": "6 人",
        "color": "C00000",
        "duty_summary": "按议程时间轴推动节目流转：主持衔接、嘉宾上下场、议程计时、颁奖执行",
        "sub_duties": [
            ("流程总控",         "持脚本对照议程；每环节超时 30s 即时提醒主持/总指挥"),
            ("舞台引导",         "嘉宾候场→上台→落座→下场；圆桌座次摆名牌"),
            ("计时员",           "主旨 25/20min、圆桌 40min；倒计时举牌（5/3/1 min）"),
            ("主持对接",         "提前与各环节主持人对脚本；过场词、串场调整"),
            ("PPT/资料对接",     "提前 24h 收齐所有讲者 PPT；按演讲顺序入备份U盘+云盘"),
            ("颁奖执行",         "奖杯/证书摆放、颁奖嘉宾引导、获奖机构合影"),
        ],
    },
    {
        "code": "D",
        "name": "媒体宣传组",
        "leader": "葛九明（组长）",
        "size": "4 人",
        "color": "ED7D31",
        "duty_summary": "覆盖现场摄影、摄像、直播、新媒体、媒体接待、会后稿件传播",
        "sub_duties": [
            ("摄影/摄像",        "主舞台机位+游机；圆桌特写；签到/茶歇/晚宴花絮"),
            ("直播/录播",        "直播平台搭建（如视频号/抖音）；导播切换；网络备用 4G"),
            ("媒体接待",         "财经/科技媒体记者签到、采访区安排、通稿派发"),
            ("新媒体即时发布",   "现场快讯（金句卡片）、嘉宾观点摘录、朋友圈/小红书推送"),
            ("会后传播",         "活动通稿、视频剪辑、嘉宾合影分发、媒体追投"),
        ],
    },
    {
        "code": "E",
        "name": "技术支持组",
        "leader": "张卓（组长）",
        "size": "3 人",
        "color": "7030A0",
        "duty_summary": "音视频、灯光、麦克风、PPT切换、网络、直播信号、应急备份",
        "sub_duties": [
            ("音响/麦克风",      "手持/领夹麦各 4 支；圆桌每位 1 支；备用电池常驻"),
            ("PPT切换/总控",     "翻页器+技术员守机；嘉宾切换提前 1 环节预加载"),
            ("灯光",             "舞台主光、圆桌追光、交响乐演出灯光预案"),
            ("LED大屏/字幕",     "主屏播放+侧屏倒计时/嘉宾介绍卡"),
            ("网络/直播",        "主有线+4G双备份；直播推流监看；断网应急脚本"),
            ("应急备件",         "备用麦×2、备用电脑×1、备用翻页器×2、转接头全套"),
        ],
    },
    {
        "code": "F",
        "name": "后勤保障组",
        "leader": "朱俊峰（组长）",
        "size": "3 人",
        "color": "548235",
        "duty_summary": "场地、物料、餐饮、茶歇、伴手礼、晚宴执行，以及陆家嘴交响乐团对接",
        "sub_duties": [
            ("场地/桌椅",        "主会场布置、圆桌摆位、签到台、媒体区、休息室"),
            ("物料制作",         "议程册、桌牌、横幅、KT板、引导牌、嘉宾胸卡"),
            ("茶歇/饮水",        "全程饮水补给；15:45音乐会期间高端茶歇（500人量）"),
            ("伴手礼",           "按等级（VVIP/嘉宾/机构）分类分发，签到台同步领取"),
            ("交响乐团对接",     "陆家嘉交响乐团到场、化妆间、乐器、调音、彩排时间"),
            ("VIP晚宴",          "席位卡、菜单、酒水、桌花、敬酒流程、餐厅对接"),
        ],
    },
    {
        "code": "G",
        "name": "应急安全/机动组",
        "leader": "刘严（组长）",
        "size": "4 人",
        "color": "BF8F00",
        "duty_summary": "安保、消防、医疗、突发事件、舆情、与场馆/公安联动",
        "sub_duties": [
            ("入场安检",         "胸卡核验、闲杂人员拒入、可疑物品筛查"),
            ("现场巡视",         "通道畅通、消防通道无堵塞、舞台周边安全"),
            ("医疗急救",         "AED 位置已知、急救包、最近医院路线（含车辆待命）"),
            ("舆情/口径",        "突发负面言论时的回应口径；与 D 组联动撤稿/澄清"),
            ("VIP/晚宴安保",     "晚宴 500 人圈层活动安全；车队动线；离场分流"),
        ],
    },
]

# =============================================================================
# 时间轴矩阵：阶段 × 各组任务
# =============================================================================
PHASES = [
    ("T-7 ~ T-3 日",  "前期筹备"),
    ("T-1 日",        "彩排与到位"),
    ("当日 09:00-13:00", "现场布置"),
    ("13:00-18:30",   "会议执行"),
    ("18:30-20:30",   "晚宴衔接"),
    ("20:30 之后",    "收尾撤场"),
]

PHASE_MATRIX = {
    # (phase_index, group_code): task
    (0,"A"): "完成总分工表、议程脚本 V1；确认嘉宾名单与到场状态",
    (0,"B"): "逐位嘉宾邀约确认；酒店/接送排期；VIP陪同名单匹配",
    (0,"C"): "议程脚本撰写；主持人对接；PPT/资料征集（D-3 截止）",
    (0,"D"): "传播节奏排期；直播平台报备；通稿初稿；机位规划",
    (0,"E"): "技术清单与场地对接；备用设备清点；网络方案",
    (0,"F"): "物料下单印刷；伴手礼采购；交响乐团合同与彩排时间确认",
    (0,"G"): "安保方案报场馆/公安备案；医疗合作机构确认；应急脚本",

    (1,"A"): "全员到岗培训；分工表过堂；对讲机分发；演练 3 次",
    (1,"B"): "VIP 接送车辆与司机对接；接机/酒店敲定；陪同清单到人",
    (1,"C"): "走台彩排：嘉宾上下场动线、圆桌座次、颁奖流程",
    (1,"D"): "直播试推流；机位实测；通稿终稿；新媒体素材预制",
    (1,"E"): "灯光/音响/PPT 全设备压力测试；备份盘 ×3 同步",
    (1,"F"): "物料到场清点；签到包/伴手礼分装；晚宴餐厅最终确认",
    (1,"G"): "安保人员到位；消防/医疗通道核查；AED 与最近医院再确认",

    (2,"A"): "现场指挥部建立；对讲机频道分配；倒计时启动",
    (2,"B"): "签到台开台；VIP 陪同上岗；接送车辆首批到位",
    (2,"C"): "主持人最终对脚本；PPT 顺序加载主机；计时器/举牌就位",
    (2,"D"): "机位/直播信号联调；签到合影区开放；记者签到",
    (2,"E"): "麦克风电池满电；灯光/音响走光走声；备用机待机",
    (2,"F"): "茶歇/水台铺设；伴手礼上签到台；交响乐团到场化妆",
    (2,"G"): "入场安检上岗；通道清空；急救点位明示",

    (3,"A"): "按议程时间轴现场指挥；超时/临时调整即时决策",
    (3,"B"): "嘉宾签到→引导→休息室→上台→落座 全程陪同",
    (3,"C"): "13:30 起按议程逐项推动；每环节剩 5/3/1 min 举牌",
    (3,"D"): "全程拍摄/直播；金句卡同步发布；记者采访区调度",
    (3,"E"): "PPT/麦克风按节目无缝切换；任何故障 30 秒内备用顶上",
    (3,"F"): "16:35 演讲间隙撤茶歇/补水；颁奖物料 18:00 前到台口",
    (3,"G"): "全场巡视；舞台/通道安全；任何医疗事件 60 秒到达",

    (4,"A"): "晚宴节奏控制；颁奖嘉宾合影；与主办方收尾交接",
    (4,"B"): "晚宴桌次引导；VIP敬酒/对接需求记录；离场用车安排",
    (4,"C"): "晚宴司仪过场词；致辞与敬酒环节衔接",
    (4,"D"): "晚宴花絮拍摄；嘉宾访谈追拍；当晚发布短视频",
    (4,"E"): "晚宴音响、背景音乐、追光；麦克风敬酒位 ×2",
    (4,"F"): "餐饮/酒水节奏；席位卡核对；伴手礼离场二次确认",
    (4,"G"): "晚宴入口二次核验；车队调度；离场分流防拥堵",

    (5,"A"): "全员复盘 30min；问题清单与改进项归档",
    (5,"B"): "嘉宾送返；当晚致谢短信/微信；伴手礼遗漏补寄",
    (5,"C"): "舞台/物料归还；脚本与签到表归档",
    (5,"D"): "通稿次日 10:00 前发布；视频 48h 内剪辑；媒体追投",
    (5,"E"): "设备清点归还；故障记录归档",
    (5,"F"): "场地交还；剩余物料盘点；账单核对",
    (5,"G"): "撤场安全确认；车队收尾；遗失物登记",
}

# =============================================================================
# 应急预案（按议程节点 + 通用类）
# =============================================================================
CONTINGENCIES = [
    # (类别, 风险, 概率, 影响, 责任组, 预案)
    ("嘉宾相关", "演讲嘉宾迟到（白硕/夏春等）",          "中", "高",
     "A+B+C",
     "B组实时位置追踪；C组备好节目顺序调换方案（先放夏春/王维军等已到场嘉宾）；A组决策；提前通知主持人衔接口径"),
    ("嘉宾相关", "嘉宾临时缺席",                          "低", "高",
     "A+C+D",
     "C组立即调整脚本，由主持人简述背景并播放主题短片；D组同步删除现场介绍物料；A组同步主办方"),
    ("嘉宾相关", "圆桌嘉宾超时发言",                     "高", "中",
     "C",
     "举牌 5/3/1min；主持人提前对话术：'最后一分钟请收口'；超 2min 主持人直接打断"),
    ("技术故障", "麦克风没声/啸叫",                       "高", "高",
     "E",
     "备用麦 30 秒内顶上；调音师切换通道；E组每个演讲前 2min 试音"),
    ("技术故障", "PPT 黑屏/翻页器失灵",                   "中", "高",
     "E+C",
     "备用电脑预加载所有 PPT；现场技术员守机手动翻页；嘉宾电脑+大会电脑双路"),
    ("技术故障", "直播断流/网络中断",                     "中", "中",
     "E+D",
     "主有线+4G 双链路自动切换；D组提前在直播间挂'信号修复中'；本地录播继续不停"),
    ("流程相关", "议程整体超时",                          "高", "高",
     "A+C",
     "音乐会前压缩茶歇/合影时间；圆桌剩 5min 时主持人直接收口；颁奖环节预留 5min 缓冲"),
    ("流程相关", "交响乐团延误",                          "低", "高",
     "F+C",
     "F组与团长保持热线；C组备用预案：将主旨3提前；提前 60min 确认到场"),
    ("人员/安全", "嘉宾或观众突发身体不适",                "中", "高",
     "G",
     "AED 与急救包就位；急救志愿者 60秒到达；车辆待命；最近医院路线已存"),
    ("人员/安全", "闲杂人员闯入/媒体冲突",                "低", "中",
     "G+D",
     "G组安保拦截；D组带记者至采访区；A组授权对外口径"),
    ("舆情",     "嘉宾发表不当言论",                       "低", "高",
     "A+D",
     "D组直播切换备用画面；A组+主办方拟统一回应口径；事后稿件删改"),
    ("后勤",     "餐饮/茶歇供应不足",                     "中", "中",
     "F",
     "茶歇按 1.2 倍备货；与餐厅签备货补货协议；F组实时巡看补给"),
    ("后勤",     "伴手礼/物料缺漏",                       "中", "低",
     "F+B",
     "签到台清单核对；缺漏当晚以快递/次日补寄"),
    ("天气",     "大雨影响嘉宾到场",                      "中", "中",
     "B+F",
     "门口加雨棚/雨伞；F组备伞 100 把；B组车辆点对点接驳"),
    ("晚宴",     "晚宴敬酒/座次冲突",                     "中", "中",
     "B+F",
     "席位卡核对到人；B组掌握 VIP 偏好；F组备用座位 2 桌"),
]

# =============================================================================
# 物资清单
# =============================================================================
# =============================================================================
# 实际志愿者分配（基于 5/13 志愿者讨论会报名意向）
# 字段：组别, 岗位, 姓名, 报名意向, 背景关键词, 分配理由
# =============================================================================
ASSIGNMENT = [
    # A 统筹协调组（3，李振春统筹工作再细分：韦佳玉接管文件秘书；李振春聚焦科创赞助）
    ("A", "总策划/总指挥",        "胡继刚（您）",      "—",         "会长/总策划",
     "全场最高决策；对接主办方与赞助方（腾讯/华为）；与春子私下推进科技展位"),
    ("A", "执行总监/副总策划+文件总管","韦佳玉",        "—（总策划点名协助统筹）","23级校友、股权一级市场投融资、参与多届",
     "扩职：原李振春的文件秘书并入；统筹议程脚本/会议记录/各组协调/复盘主持"),
    ("A", "科创资源对接专员",     "李振春（春子）",     "未明示",     "经院校友、科创业务、有会务/签到经验",
     "★职责再细分：聚焦科创赞助谈判+科技展位推进（自己提出的科技展示建议由他与总策划主推）"),

    # B 嘉宾接待组（5）—— Winny 调出至 F
    ("B", "组长",                 "王胜（20级）",      "B/F → B",   "第三次参加最资深",
     "总策划指定升任：第三次参加经验最足，带 B 组"),
    ("B", "VIP 1V1（港澳/海外嘉宾）","朱铭喆",          "B/C → B",   "友邦保险、香港城市大学上海联络人",
     "港校联络人，专业形象+沟通能力，适合涉外/港澳嘉宾"),
    ("B", "VIP 1V1（主旨嘉宾）",   "JasonCAI 蔡杰",     "B/G → B",   "复旦MBA、设计咨询",
     "形象气质+审美佳，适合白硕、夏春等主旨嘉宾"),
    ("B", "观众席引导",           "Angi 郭杰",         "B",         "国际贸易、首次",
     "采纳冯墨建议补观众席引导盲区，新人友好岗"),
    ("B", "动线引导",             "LV",               "B",         "—",
     "标准引导岗"),

    # C 会务执行组（6）—— + 吕志翔（从 E 调入）
    ("C", "组长 / 舞台+设备总管", "陈潇 Kelly",        "C",         "广告活动一条龙、张江有自有灯光/舞台/音响工厂、核心搭子",
     "总策划指定升任：20年策划经验+自有设备，全场制作把关"),
    ("C", "副组长 / 计时员",      "王珏",              "C",         "中信证券，自报会务+倒计时",
     "保留计时核心岗（自荐），副组长稳定流程节奏"),
    ("C", "主持对接/串场",        "冯墨 律师",         "C/B → C",   "私募→律师、时间灵活、沟通强",
     "律师沟通力+时间灵活，最适合协调主持脚本"),
    ("C", "PPT/资料对接",         "黄璐 Lucius",       "B/C → C",   "建筑设计+地产+养老研究、首次",
     "细致岗，提前 24h 收齐讲者 PPT"),
    ("C", "颁奖执行/机构对接",    "徐胜博",            "未明示",     "一二级创投、K12教育",
     "创投背景对接颁奖机构得体"),
    ("C", "圆桌/机构嘉宾对接（调入）","吕志翔",         "B→E→C",     "12级、投融资、教育/医疗/AI",
     "投融资人脉，对接 4 位圆桌嘉宾+主旨讲者上下场（技术非其专长，回归会务）"),

    # D 媒体宣传组（4）—— 马磊调出至 G
    ("D", "组长",                "葛九明",            "—（总策划指定）","媒体/宣传背景（待补全）",
     "总策划指定 D 组组长，统筹摄影摄像/直播/新媒体/媒体接待对外口径"),
    ("D", "副组长 / 摄影摄像+自媒体","张蒙 AEGISTAR",  "D",         "大宗行/AI数据中心、本身做自媒体",
     "保留摄影摄像与自媒体核心岗，作为 D 组副组长"),
    ("D", "AI 内容/新媒体",       "蒋珊",             "B → D",     "AI赛道、复旦人工智能俱乐部、跨校 C9 科创金融研究院",
     "AI 圈层资源最强；新媒体（金句卡/朋友圈/小红书）"),
    ("D", "直播+交响乐预热+媒体接待","皮尔德小号",     "未明示",     "互联网 10 年→城市更新、资源对接强",
     "扩职：原马磊的媒体接待并入；主管直播+交响乐往届视频 T-7 起视频号每日预热"),

    # E 技术支持组（3）—— 张卓任组长（律师做协调）；吕志翔调出至 C
    ("E", "组长 / 协调统筹（调入）","张卓 盈科",       "B/G → E",   "财税律师、二次参加、协调能力强",
     "总策划指定：律师协调能力强，统筹 E 组与场馆 AV 外包/陈潇工厂设备"),
    ("E", "副组长 / 技术执行",     "随圣博 Joshua-Sui","B/E/G → E", "复旦管院年会后勤经验",
     "技术执行核心，对接场馆 AV 现场操作"),
    ("E", "技术对接助理",         "高辰辰",            "B → E（听安排）","同浦汇、首次",
     "标准化对接执行"),

    # F 后勤保障组（3）—— 刘严调走任 G 组长；Winny 调入做副组长
    ("F", "组长 / 场地+物料+赞助物料","朱俊峰",         "F",         "19级、建筑+市场开发、复旦网球协会运营",
     "扩职：原刘严的赞助物料对接并入；建筑+协会运营经验主导场地与物料"),
    ("F", "副组长 / 茶歇+晚宴主管（调入）","Winny 温妮","B/C/F均可 → F","空窗期时间最充裕、IP授权+MINISO运营",
     "★最适合茶歇/晚宴：空窗期时间最多+零售运营+商务接待经验；从 B 调入做副组长"),
    ("F", "茶歇+晚宴专员",        "彭常丽",            "B/F → F",   "金融、二次参与",
     "去年已熟悉流程，副手 Winny 主管 500 人茶歇+晚宴对接"),

    # G 应急/机动组（4）—— 刘严任组长，蔡萍升副组长，马磊调入
    ("G", "组长 / 安保+机动调配（调入）","刘严",       "F → G",     "自报积极、平日接活多、协调力强",
     "总策划指定：从 F 副组长调任 G 组长，统筹安保+机动跨组调配"),
    ("G", "副组长 / 急救+医疗主岗 ★","蔡萍",           "B/G → G",   "复旦经院培训负责人(退休)、★持有急救证",
     "升任副组长；保留急救/医疗专项岗（持证稀缺技能）"),
    ("G", "机动 / 媒体应急口径（调入）","马磊",         "D → G",     "AI Club 生态、跨界资源",
     "媒体应急口径+跨界资源；G 组机动支援"),
    ("G", "机动 / 安保",          "康诺斯 Conns",      "B/G → G",   "设计改造+碳资产、首次",
     "保留 G 组（自选偏好），机动支援+入场安检"),
]

# =============================================================================
# 核心策划圈（总策划直接联动的小圈子，跨组高频沟通）
# =============================================================================
CORE_CIRCLE = [
    ("胡继刚（您）",  "总策划",        "全场指挥"),
    ("韦佳玉",       "A 执行总监",    "副总策划，统筹各组+文件总管"),
    ("陈潇 Kelly",   "C 组组长",     "提供灯光/舞台/音响设备资源"),
    ("王胜",         "B 组组长",     "三次参加，老带新"),
    ("葛九明",       "D 组组长",     "总策划指定，统筹媒体宣传"),
    ("刘严",         "G 组组长",     "应急安保+机动跨组调配"),
    ("李振春 春子", "A 组科创对接", "科创赞助/科技展位推进"),
]

# =============================================================================
# 5/13 优化建议筛选清单（采纳/调整/暂缓）
# =============================================================================
OPT_SUGGESTIONS = [
    # (建议人, 建议内容, 处置, 责任组, 落地动作)
    ("韦佳玉", "量化各组任务量+设置 Milestone（按时间精力选岗位）",
     "采纳", "A",   "新增 Milestone 表（见 09 sheet），各组组长 T-7 前提交"),
    ("韦佳玉", "所有 Milestone 节点提前 1-2 天，留 buffer",
     "采纳", "A",   "阶段任务矩阵已整体前移，T-3 改为 T-5"),
    ("张卓",   "提前半天/线上彩排，明确工作范围边界",
     "采纳", "A+全员","T-1 18:00 全员线上彩排 1.5h；T-7 现场踏勘 1 次"),
    ("张卓",   "协助人员对接清单（酒店/屏幕谁联系，避免事事请示总策划）",
     "采纳", "A",   "新增 11-对接联系人清单（见 sheet），人手一份"),
    ("张卓",   "嘉宾发言时间提前打招呼+主持人收口",
     "采纳", "C",   "C 组 T-3 致电每位讲者+主持人对脚本（见 10-嘉宾时间SOP）"),
    ("春子",   "增加科技类展示（AI/机器人/芯片小海报）",
     "暂缓→单列待办", "A",
     "需赞助谈判，春子+胡继刚 私下推进，不进 5/22 必交付"),
    ("蒋珊",   "扩大交响乐宣传，视频号发去年录像预告",
     "采纳", "D",   "D 组皮尔德小号负责，T-7 起每日 1 条；T-3 集中投流"),
    ("蒋珊",   "演奏者走位与观众互动",
     "暂缓",  "—", "总策划已明确专业度不足；现场仅做灯光烘托"),
    ("隋胜博/徐胜博","地铁→场地快进短视频群发",
     "采纳", "D+F","T-5 拍摄、T-3 发群（去年已做，今年延续）"),
    ("JasonCAI","T-2 微信+短信点对点提醒嘉宾交通/停车",
     "采纳", "B",   "B 组 VIP 1V1 责任人 T-2 18:00 前完成全部联络"),
    ("王胜",   "组长 A/B 角备份制度",
     "采纳", "A",   "新增 12-组长备份表（见 sheet）每组指定 B 角"),
    ("康诺斯", "B 组按身高/年龄筛选",
     "调整后采纳", "B",
     "改为按英语/沟通/形象气质匹配 VIP 1V1，避免歧视"),
    ("张蒙",   "志愿者专用颜色胸牌快速识别",
     "采纳", "F",   "F 组制作橙色胸牌（与嘉宾胸卡明显区分），见物资清单"),
    ("冯墨",   "现场观众管理（盲区）",
     "采纳", "B+G","B 组 Angi 改为观众席引导岗；G 组机动支援"),
    ("蔡萍",   "自报有急救证",
     "采纳", "G",   "蔡萍调到 G 组任急救/医疗主岗（稀缺技能）"),
    ("陈潇",   "张江工厂提供灯光/舞台/音响设备",
     "采纳", "C+E","作为重大物料资源并入；C 组陈潇为设备总管，E 组对接"),
    ("彭常丽", "去年赞助物料环节混乱，要提前深化",
     "采纳", "F",   "F 组刘严专项负责赞助物料入场+点位"),
    ("总策划确认","工作装+合影+证书+午餐+纪念品",
     "采纳", "A+F","已写入物资清单；T-1 全员告知"),
    ("总策划确认","核心搭子：韦佳玉/Winnie/陈潇/王胜+葛九明",
     "采纳", "A",   "建立核心策划圈高频沟通群（含 D 组长葛九明 7 人）"),
]

# =============================================================================
# Milestone（节点提前 1-2 天）
# =============================================================================
MILESTONES = [
    # (时间节点, 关键里程碑, 责任组)
    ("T-9 (5/13)",  "志愿者讨论会完成→分工表发布（已完成）",        "A"),
    ("T-8 (5/14)",  "各组建独立沟通群+组长 A/B 角确定+核心策划圈群","A+各组长"),
    ("T-8 (5/14)",  "组长见面会（线上 1h）",                         "A+组长"),
    ("T-7 (5/15)",  "晚宴名单锁定 V1（去年痛点：太晚锁定）",          "B+F"),
    ("T-7 (5/15)",  "交响乐往届视频开始日更预热（视频号）",          "D"),
    ("T-7 (5/15)",  "现场踏勘 1（场地/动线/AV）",                    "A+C+E+F"),
    ("T-5 (5/17)",  "所有讲者 PPT 截稿（去年是 T-3，今年提前）",      "C"),
    ("T-5 (5/17)",  "地铁→场地快进短视频拍摄完成",                   "D+F"),
    ("T-5 (5/17)",  "物料/胸牌印刷下单（去年痛点：名牌制作晚）",      "F"),
    ("T-4 (5/18)",  "腾讯/华为赞助物料对接清单确认",                 "F+A"),
    ("T-3 (5/19)",  "C 组致电每位讲者+主持人对发言时间（防超时）",    "C"),
    ("T-3 (5/19)",  "对接联系人清单（酒店/AV/餐饮/安保）发到全员",    "A"),
    ("T-3 (5/19)",  "短视频发主群+嘉宾群预热",                       "D"),
    ("T-2 (5/20)",  "B 组 VIP 1V1 点对点短信/微信提醒嘉宾",          "B"),
    ("T-2 (5/20)",  "现场踏勘 2（含彩排站位）",                       "A+C+E"),
    ("T-1 (5/21)",  "全员线上彩排 18:00-19:30（1.5h）",              "A+全员"),
    ("T-1 (5/21)",  "全设备压力测试+陈潇工厂设备进场",                "C+E"),
    ("T-0 09:00",   "现场指挥部成立、对讲机分发、AED/急救包就位",     "A+G"),
    ("T-0 12:00",   "签到台开台、VIP 车队首班、直播试推流",           "B+D"),
    ("T-0 13:25",   "全员就位静默 5 分钟、广播开场倒计时",            "A"),
    ("T+1",         "通稿次日 10:00 前发布、复盘会、问题归档",        "A+D"),
]

# =============================================================================
# 嘉宾发言时间控制 SOP（去年最大事故）
# =============================================================================
TIME_SOP = [
    ("T-7", "整理每位讲者的演讲时长合同/邀约函（25/20/40min）", "C"),
    ("T-3", "C 组致电每位讲者：再次确认时长 + 提示我们会举牌",   "C"),
    ("T-3", "C 组与 4 位主持人对脚本：嘉宾超时由你直接打断的话术",  "C"),
    ("T-1", "彩排时给主持人发『收口话术卡』：'感谢XXX，最后一分钟'", "C"),
    ("T-0 演讲前 2min", "计时员就位，倒计时牌(5/3/1 min)摆放台口",   "C"),
    ("T-0 演讲中 5/3/1min", "在嘉宾视线方向举牌（先黄后红）",         "C"),
    ("T-0 超时 30s",   "主持人通过耳麦提示'最后一句话'",                  "C"),
    ("T-0 超时 90s",   "主持人直接收口：'感谢XXX的精彩分享'走流程",      "C"),
    ("T-0 圆桌嘉宾",   "每人剩 1min 时主持人示意；超时立即转下一位",       "C"),
]

# =============================================================================
# 协助人员对接清单（解决"事事请示总策划"瓶颈）
# =============================================================================
LIAISON_LIST = [
    # (问题类别, 找谁, 备用)
    ("场地/桌椅/赞助物料入场",   "F 组 朱俊峰",       "韦佳玉（A 组）"),
    ("灯光/舞台/音响（自有）",   "C 组 陈潇 Kelly",   "E 组 张卓"),
    ("麦克风/PPT切换/网络/直播",  "E 组 随圣博",       "E 组 高辰辰"),
    ("E 组协调/对接外包 AV",      "E 组 张卓（组长）",  "随圣博（B 角）"),
    ("酒店/嘉宾接送/休息室",    "B 组 王胜",         "韦佳玉（A 组）"),
    ("VIP 嘉宾临时需求",        "B 组对应 1V1 责任人","王胜（B 组长）"),
    ("观众席引导/秩序",          "B 组 Angi",        "G 组 康诺斯"),
    ("茶歇 / 晚宴 ★",           "F 组 Winny（茶歇晚宴主管）","彭常丽 / 朱俊峰"),
    ("摄影/摄像",               "D 组 张蒙",         "葛九明（D 组长）"),
    ("直播信号/新媒体即时发布",  "D 组 皮尔德小号",   "蒋珊（AI 内容）"),
    ("AI 圈层金句/新媒体素材",   "D 组 蒋珊",         "皮尔德小号"),
    ("媒体/记者接待",           "D 组 皮尔德小号",   "葛九明"),
    ("讲者 PPT/资料/换片",      "C 组 黄璐",         "陈潇 Kelly"),
    ("流程超时/嘉宾迟到",       "C 组 王珏（计时）",  "陈潇 → 韦佳玉"),
    ("圆桌嘉宾上下场",          "C 组 吕志翔",       "王珏"),
    ("颁奖物料/获奖嘉宾",       "C 组 徐胜博",       "王珏"),
    ("法务/舆情/突发声明",       "C 组 冯墨（律师）",  "E 组 张卓（律师）→ 胡继刚"),
    ("医疗/急救 ★",             "G 组 蔡萍 ★",       "G 组 刘严"),
    ("跨组人手不足/机动调配",    "G 组 刘严（组长）",  "G 组 康诺斯 / 马磊"),
    ("媒体应急口径",             "G 组 马磊",        "D 组 葛九明"),
    ("入场安检/可疑人员",         "G 组 康诺斯",       "G 组 刘严"),
    ("交响乐团对接",             "F 组 朱俊峰",      "韦佳玉（去年负责）"),
    ("晚宴席位/敬酒/现场对接",   "F 组 Winny",       "B 组 王胜"),
    ("科创展位/赞助谈判",        "A 组 春子",         "胡继刚"),
    ("议程脚本/会议记录/复盘",   "A 组 韦佳玉",       "胡继刚"),
    ("无法解决/重大决策",        "韦佳玉 → 胡继刚",  "—"),
]

# =============================================================================
# 组长 A/B 角备份制度（应对组长当天缺席）
# =============================================================================
BACKUP_ROLES = [
    # (组, A 角组长, B 角备份, 备注)
    ("A 统筹",   "胡继刚（总策划）", "韦佳玉",        "总策划缺席由韦佳玉全权代行；春子专注科创赞助"),
    ("B 接待",   "王胜",          "朱铭喆",        "A 角第三次参加最资深；B 角港校联络人"),
    ("C 会务",   "陈潇 Kelly",    "王珏",          "A 角 20 年策划经验+自有设备；B 角中信证券计时主控"),
    ("D 媒体",   "葛九明",        "张蒙 AEGISTAR", "A 角总策划指定；B 角自媒体+摄影核心岗"),
    ("E 技术",   "张卓 盈科",     "随圣博",        "A 角律师协调统筹；B 角技术执行（操作设备）"),
    ("F 后勤",   "朱俊峰",        "Winny 温妮",    "A 角建筑+协会运营；B 角空窗期时间最多+零售运营"),
    ("G 应急",   "刘严",          "蔡萍",          "A 角统筹机动调配；B 角持急救证常驻医疗岗★"),
]

MATERIALS = [
    ("会务", "议程手册",                      "600 份",                "F",   "T-5 印刷完成（前移）"),
    ("会务", "桌牌/嘉宾胸卡",                 "全员",                  "F",   "T-2 完成（去年痛点：太晚）"),
    ("会务", "★志愿者专用橙色胸牌",           "26 张（含韦佳玉）",     "F",   "T-2 完成（采纳张蒙建议）"),
    ("会务", "签到表/电子签到",               "分级 4 类",             "F+B", "当日 12:00 上线"),
    ("会务", "横幅/KT板/引导牌",              "全场",                  "F",   "当日 10:00 前"),
    ("技术", "手持麦×4 / 领夹麦×4",            "各 2 备用",             "E",   "当日 10:00 测试"),
    ("技术", "翻页器×4 / 备用电脑×1",          "——",                    "E",   "T-1 检查"),
    ("技术", "转接头/线材",                   "HDMI/Type-C/Lightning 各 3","E","T-1"),
    ("技术", "对讲机",                        "26 台（分 3 频道）",    "A+E", "当日 09:00 分发"),
    ("技术", "★陈潇张江工厂灯光/舞台/音响",   "1 套（自有资源）",      "C+E", "T-1 进场（陈潇主导，节省赞助）"),
    ("媒体", "摄影/摄像机",                   "主机×2 游机×2",         "D",   "当日 11:00 调试"),
    ("媒体", "直播推流设备",                  "1 套（含备用 4G）",     "D+E", "T-1 试推（直播确认开启）"),
    ("媒体", "★地铁→场地快进短视频",          "1 段 60 秒",           "D+F", "T-5 拍摄完成、T-3 群发"),
    ("媒体", "★交响乐往届预热视频",           "5-7 段 短视频",         "D",   "T-7 起视频号每日 1 条"),
    ("后勤", "茶歇（500人）",                 "高端冷餐/茶饮",         "F",   "15:30 上桌"),
    ("后勤", "伴手礼",                        "VVIP/嘉宾/机构 3 档",   "F",   "T-2 分装（前移）"),
    ("后勤", "晚宴桌花/席位卡",               "500 人 50 桌",          "F",   "18:00 前完成（晚宴名单 T-7 锁定）"),
    ("后勤", "★赞助物料（腾讯/华为）",        "按合同清单",            "F",   "T-4 与品牌方对接确认（去年痛点）"),
    ("后勤", "志愿者午餐+证书+纪念品",        "26 套",                "A+F", "T-1 分装、当日 12:00 发"),
    ("应急", "急救包/AED 标识",               "2 套 / 全场 3 处",      "G",   "当日 09:00 就位"),
    ("应急", "★急救主岗（蔡萍持证）",         "1 名",                 "G",   "全程在 G 组指挥点待命"),
    ("应急", "应急联络卡（含对接清单）",      "全员人手 1 张",         "A+G", "T-3 发放（采纳张卓建议）"),
    ("应急", "备用雨伞",                      "100 把",               "F",   "视天气"),
    ("通知", "★T-2 嘉宾点对点提醒短信/微信",  "全部嘉宾",             "B",   "T-2 18:00 前完成（采纳JasonCAI建议）"),
]

# =============================================================================
# 工具函数
# =============================================================================
def hex_to_rgb(h):
    h = h.lstrip("#")
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

# =============================================================================
# 生成 Excel
# =============================================================================
def build_excel(path):
    wb = Workbook()

    thin = Side(border_style="thin", color="BFBFBF")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    header_fill = PatternFill("solid", fgColor="1F4E79")
    header_font = Font(name="微软雅黑", size=11, bold=True, color="FFFFFF")
    cell_font   = Font(name="微软雅黑", size=10)
    title_font  = Font(name="微软雅黑", size=16, bold=True, color="1F4E79")
    center      = Alignment(horizontal="center", vertical="center", wrap_text=True)
    left        = Alignment(horizontal="left", vertical="center", wrap_text=True)

    def style_header(ws, row, ncols):
        for c in range(1, ncols+1):
            cell = ws.cell(row=row, column=c)
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = center
            cell.border = border

    def style_body(ws, r1, r2, ncols, align=left):
        for r in range(r1, r2+1):
            for c in range(1, ncols+1):
                cell = ws.cell(row=r, column=c)
                cell.font = cell_font
                cell.alignment = align
                cell.border = border

    # ----- Sheet1 总览 -----
    ws = wb.active
    ws.title = "01-总览"
    ws["A1"] = EVENT_TITLE
    ws["A1"].font = title_font
    ws.merge_cells("A1:F1")
    ws["A2"] = f"愿景：{EVENT_SUB}"
    ws["A3"] = f"时间：{EVENT_DATE}"
    ws["A4"] = f"地点：{EVENT_VENUE}"
    for r in (2,3,4):
        ws.cell(row=r, column=1).font = Font(name="微软雅黑", size=11, bold=True, color="404040")
        ws.merge_cells(start_row=r, start_column=1, end_row=r, end_column=6)

    ws["A6"] = "志愿者总组织架构（向总策划负责）"
    ws["A6"].font = Font(name="微软雅黑", size=13, bold=True, color="C00000")
    ws.merge_cells("A6:F6")

    headers = ["组别代码","组别","负责人","建议人数","主要职责","对接组别"]
    for i, h in enumerate(headers, 1):
        ws.cell(row=7, column=i, value=h)
    style_header(ws, 7, len(headers))

    for i, g in enumerate(GROUPS):
        r = 8 + i
        ws.cell(row=r, column=1, value=g["code"])
        ws.cell(row=r, column=2, value=g["name"])
        ws.cell(row=r, column=3, value=g["leader"])
        ws.cell(row=r, column=4, value=g["size"])
        ws.cell(row=r, column=5, value=g["duty_summary"])
        ws.cell(row=r, column=6, value="全员对接 A 组")
    style_body(ws, 8, 8+len(GROUPS)-1, len(headers))

    widths = [10, 16, 22, 12, 60, 16]
    for i, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = w
    for r in range(8, 8+len(GROUPS)):
        ws.row_dimensions[r].height = 42

    # ----- Sheet2 各组细分职责 -----
    ws = wb.create_sheet("02-各组细分职责")
    ws["A1"] = "各组职能细分（大项 → 小项）"
    ws["A1"].font = title_font
    ws.merge_cells("A1:D1")

    headers = ["组别","子职能/岗位","具体职责","建议负责志愿者"]
    row = 3
    for i, h in enumerate(headers, 1):
        ws.cell(row=row, column=i, value=h)
    style_header(ws, row, len(headers))
    row += 1
    body_start = row
    for g in GROUPS:
        group_label = f"{g['code']} {g['name']}"
        start = row
        for sub_name, sub_desc in g["sub_duties"]:
            ws.cell(row=row, column=1, value=group_label)
            ws.cell(row=row, column=2, value=sub_name)
            ws.cell(row=row, column=3, value=sub_desc)
            ws.cell(row=row, column=4, value="")
            row += 1
        if row - 1 > start:
            ws.merge_cells(start_row=start, start_column=1, end_row=row-1, end_column=1)
    body_end = row - 1
    style_body(ws, body_start, body_end, len(headers))
    for r in range(body_start, body_end+1):
        ws.row_dimensions[r].height = 32
    widths = [18, 22, 70, 20]
    for i, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = w

    # ----- Sheet3 阶段 × 组别 任务矩阵 -----
    ws = wb.create_sheet("03-阶段任务矩阵")
    ws["A1"] = "阶段 × 组别 任务矩阵（T-7 至 收尾）"
    ws["A1"].font = title_font
    ws.merge_cells("A1:I1")

    headers = ["阶段","时段"] + [f"{g['code']} {g['name']}" for g in GROUPS]
    row = 3
    for i, h in enumerate(headers, 1):
        ws.cell(row=row, column=i, value=h)
    style_header(ws, row, len(headers))

    row += 1
    body_start = row
    for pi, (when, name) in enumerate(PHASES):
        ws.cell(row=row, column=1, value=name)
        ws.cell(row=row, column=2, value=when)
        for gi, g in enumerate(GROUPS):
            ws.cell(row=row, column=3+gi, value=PHASE_MATRIX.get((pi, g["code"]), ""))
        row += 1
    body_end = row - 1
    style_body(ws, body_start, body_end, len(headers))
    for r in range(body_start, body_end+1):
        ws.row_dimensions[r].height = 70
    widths = [14, 16, 32, 32, 32, 32, 32, 32, 32]
    for i, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = w

    # ----- Sheet4 议程时间轴分工 -----
    ws = wb.create_sheet("04-议程时间轴分工")
    ws["A1"] = "议程时间轴 × 责任组 × 关键动作"
    ws["A1"].font = title_font
    ws.merge_cells("A1:D1")

    headers = ["时间","议程节点","主责组","关键动作/物料"]
    row = 3
    for i, h in enumerate(headers, 1):
        ws.cell(row=row, column=i, value=h)
    style_header(ws, row, len(headers))

    row += 1
    body_start = row
    for t, name, groups, action in AGENDA:
        ws.cell(row=row, column=1, value=t)
        ws.cell(row=row, column=2, value=name)
        ws.cell(row=row, column=3, value=groups)
        ws.cell(row=row, column=4, value=action)
        row += 1
    body_end = row - 1
    style_body(ws, body_start, body_end, len(headers))
    for r in range(body_start, body_end+1):
        ws.row_dimensions[r].height = 28
    widths = [16, 38, 14, 60]
    for i, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = w

    # ----- Sheet5 应急预案 -----
    ws = wb.create_sheet("05-应急预案")
    ws["A1"] = "应急预案 - 风险类型 × 责任组 × 处置流程"
    ws["A1"].font = title_font
    ws.merge_cells("A1:F1")

    headers = ["类别","风险事件","概率","影响","责任组","处置预案"]
    row = 3
    for i, h in enumerate(headers, 1):
        ws.cell(row=row, column=i, value=h)
    style_header(ws, row, len(headers))

    row += 1
    body_start = row
    for cat, risk, p, imp, resp, plan in CONTINGENCIES:
        ws.cell(row=row, column=1, value=cat)
        ws.cell(row=row, column=2, value=risk)
        ws.cell(row=row, column=3, value=p)
        ws.cell(row=row, column=4, value=imp)
        ws.cell(row=row, column=5, value=resp)
        ws.cell(row=row, column=6, value=plan)
        row += 1
    body_end = row - 1
    style_body(ws, body_start, body_end, len(headers))
    for r in range(body_start, body_end+1):
        ws.row_dimensions[r].height = 40
    widths = [12, 28, 8, 8, 14, 60]
    for i, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = w

    # ----- Sheet6 物资清单 -----
    ws = wb.create_sheet("06-物资清单")
    ws["A1"] = "物资 / 设备 清单"
    ws["A1"].font = title_font
    ws.merge_cells("A1:E1")

    headers = ["类别","物品","数量","负责组","到位时间"]
    row = 3
    for i, h in enumerate(headers, 1):
        ws.cell(row=row, column=i, value=h)
    style_header(ws, row, len(headers))

    row += 1
    body_start = row
    for cat, item, qty, resp, when in MATERIALS:
        ws.cell(row=row, column=1, value=cat)
        ws.cell(row=row, column=2, value=item)
        ws.cell(row=row, column=3, value=qty)
        ws.cell(row=row, column=4, value=resp)
        ws.cell(row=row, column=5, value=when)
        row += 1
    body_end = row - 1
    style_body(ws, body_start, body_end, len(headers))
    for r in range(body_start, body_end+1):
        ws.row_dimensions[r].height = 24
    widths = [10, 30, 28, 14, 18]
    for i, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = w

    # ----- Sheet7 实际志愿者分配（基于 5/13 讨论会） -----
    ws = wb.create_sheet("07-志愿者实际分配")
    ws["A1"] = "志愿者实际分配（基于 2026-05-13 志愿者讨论会报名意向）"
    ws["A1"].font = title_font
    ws.merge_cells("A1:F1")
    ws["A2"] = "共 25 人；偏好→实分 表示该志愿者勾选了多组、最终分配的组别"
    ws["A2"].font = Font(name="微软雅黑", size=10, color="808080")
    ws.merge_cells("A2:F2")

    headers = ["组别","岗位","姓名","报名意向","背景关键词","分配理由"]
    row = 4
    for i, h in enumerate(headers, 1):
        ws.cell(row=row, column=i, value=h)
    style_header(ws, row, len(headers))

    group_colors = {g["code"]: g["color"] for g in GROUPS}

    row += 1
    start = row
    for code, post, name, pref, bg, reason in ASSIGNMENT:
        ws.cell(row=row, column=1, value=code)
        ws.cell(row=row, column=2, value=post)
        ws.cell(row=row, column=3, value=name)
        ws.cell(row=row, column=4, value=pref)
        ws.cell(row=row, column=5, value=bg)
        ws.cell(row=row, column=6, value=reason)
        # group code cell colored
        c = ws.cell(row=row, column=1)
        c.fill = PatternFill("solid", fgColor=group_colors.get(code, "808080"))
        c.font = Font(name="微软雅黑", size=11, bold=True, color="FFFFFF")
        c.alignment = center
        row += 1
    end = row - 1
    style_body(ws, start, end, len(headers))
    # restyle group code column
    for r in range(start, end+1):
        c = ws.cell(row=r, column=1)
        c.alignment = center
        ws.row_dimensions[r].height = 34
    widths = [8, 26, 20, 16, 38, 50]
    for i, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = w

    # ----- Sheet8 联络方式空模板（请填） -----
    ws = wb.create_sheet("08-志愿者联络表")
    ws["A1"] = "志愿者联络方式（请补全）"
    ws["A1"].font = title_font
    ws.merge_cells("A1:G1")

    headers = ["组别","姓名","手机","微信","到场时间","应急联系人","备注"]
    row = 3
    for i, h in enumerate(headers, 1):
        ws.cell(row=row, column=i, value=h)
    style_header(ws, row, len(headers))

    row += 1
    start = row
    for code, post, name, *_ in ASSIGNMENT:
        ws.cell(row=row, column=1, value=f"{code} - {post}")
        ws.cell(row=row, column=2, value=name)
        row += 1
    end = row - 1
    style_body(ws, start, end, len(headers))
    widths = [22, 18, 16, 16, 18, 18, 24]
    for i, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = w

    # ----- Sheet9 5/13 优化建议筛选清单 -----
    ws = wb.create_sheet("09-优化建议筛选")
    ws["A1"] = "5/13 志愿者讨论会  19 条建议筛选与处置"
    ws["A1"].font = title_font
    ws.merge_cells("A1:E1")
    ws["A2"] = "处置：采纳 / 调整后采纳 / 暂缓（暂缓 = 经判断不进 5/22 必交付，但不否定其价值）"
    ws["A2"].font = Font(name="微软雅黑", size=10, color="808080")
    ws.merge_cells("A2:E2")

    headers = ["建议人", "建议内容", "处置", "责任组", "落地动作"]
    row = 4
    for i, h in enumerate(headers, 1):
        ws.cell(row=row, column=i, value=h)
    style_header(ws, row, len(headers))

    row += 1
    start = row
    status_color = {
        "采纳": "548235",
        "调整后采纳": "ED7D31",
        "暂缓": "808080",
        "暂缓→单列待办": "808080",
    }
    for sug, content, status, resp, act in OPT_SUGGESTIONS:
        ws.cell(row=row, column=1, value=sug)
        ws.cell(row=row, column=2, value=content)
        ws.cell(row=row, column=3, value=status)
        ws.cell(row=row, column=4, value=resp)
        ws.cell(row=row, column=5, value=act)
        c = ws.cell(row=row, column=3)
        c.fill = PatternFill("solid", fgColor=status_color.get(status, "808080"))
        c.font = Font(name="微软雅黑", size=10, bold=True, color="FFFFFF")
        c.alignment = center
        row += 1
    end = row - 1
    style_body(ws, start, end, len(headers))
    for r in range(start, end+1):
        ws.cell(row=r, column=3).alignment = center
        ws.row_dimensions[r].height = 36
    widths = [12, 50, 16, 14, 50]
    for i, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = w

    # ----- Sheet10 Milestone 时间表 -----
    ws = wb.create_sheet("10-Milestone时间表")
    ws["A1"] = "Milestone 时间表（采纳韦佳玉建议：所有节点提前 1-2 天）"
    ws["A1"].font = title_font
    ws.merge_cells("A1:C1")

    headers = ["时间节点", "关键里程碑", "责任组"]
    row = 3
    for i, h in enumerate(headers, 1):
        ws.cell(row=row, column=i, value=h)
    style_header(ws, row, len(headers))

    row += 1
    start = row
    for t, ms, resp in MILESTONES:
        ws.cell(row=row, column=1, value=t)
        ws.cell(row=row, column=2, value=ms)
        ws.cell(row=row, column=3, value=resp)
        row += 1
    end = row - 1
    style_body(ws, start, end, len(headers))
    for r in range(start, end+1):
        ws.row_dimensions[r].height = 28
    widths = [16, 60, 16]
    for i, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = w

    # ----- Sheet11 嘉宾发言时间控制 SOP -----
    ws = wb.create_sheet("11-嘉宾时间控制SOP")
    ws["A1"] = "嘉宾发言时间控制 SOP（去年最大事故：从第 1 位讲者就开始超时）"
    ws["A1"].font = title_font
    ws.merge_cells("A1:C1")

    headers = ["时间点", "动作", "责任组"]
    row = 3
    for i, h in enumerate(headers, 1):
        ws.cell(row=row, column=i, value=h)
    style_header(ws, row, len(headers))

    row += 1
    start = row
    for t, act, resp in TIME_SOP:
        ws.cell(row=row, column=1, value=t)
        ws.cell(row=row, column=2, value=act)
        ws.cell(row=row, column=3, value=resp)
        row += 1
    end = row - 1
    style_body(ws, start, end, len(headers))
    for r in range(start, end+1):
        ws.row_dimensions[r].height = 28
    widths = [22, 60, 14]
    for i, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = w

    # ----- Sheet12 协助人员对接清单 -----
    ws = wb.create_sheet("12-对接联系人清单")
    ws["A1"] = "协助人员对接清单（解决『事事请示总策划』瓶颈，全员人手一份）"
    ws["A1"].font = title_font
    ws.merge_cells("A1:C1")
    ws["A2"] = "采纳张卓建议：现场遇到问题先找对应志愿者，找不到再找 B 角，最后才找韦佳玉/胡继刚"
    ws["A2"].font = Font(name="微软雅黑", size=10, color="808080")
    ws.merge_cells("A2:C2")

    headers = ["问题类别 / 你遇到的事情", "找谁（A 角）", "找不到时（B 角）"]
    row = 4
    for i, h in enumerate(headers, 1):
        ws.cell(row=row, column=i, value=h)
    style_header(ws, row, len(headers))

    row += 1
    start = row
    for cat, who, backup in LIAISON_LIST:
        ws.cell(row=row, column=1, value=cat)
        ws.cell(row=row, column=2, value=who)
        ws.cell(row=row, column=3, value=backup)
        row += 1
    end = row - 1
    style_body(ws, start, end, len(headers))
    for r in range(start, end+1):
        ws.row_dimensions[r].height = 26
    widths = [42, 28, 28]
    for i, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = w

    # ----- Sheet13 组长 A/B 角备份制度 -----
    ws = wb.create_sheet("13-组长AB角备份")
    ws["A1"] = "组长 A/B 角备份制度（采纳王胜建议：应对组长当天缺席）"
    ws["A1"].font = title_font
    ws.merge_cells("A1:D1")

    headers = ["组", "A 角组长", "B 角备份", "备注"]
    row = 3
    for i, h in enumerate(headers, 1):
        ws.cell(row=row, column=i, value=h)
    style_header(ws, row, len(headers))

    row += 1
    start = row
    for grp, a, b, note in BACKUP_ROLES:
        ws.cell(row=row, column=1, value=grp)
        ws.cell(row=row, column=2, value=a)
        ws.cell(row=row, column=3, value=b)
        ws.cell(row=row, column=4, value=note)
        row += 1
    end = row - 1
    style_body(ws, start, end, len(headers))
    for r in range(start, end+1):
        ws.row_dimensions[r].height = 28
    widths = [16, 20, 20, 50]
    for i, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = w

    # ----- Sheet14 核心策划圈 -----
    ws = wb.create_sheet("14-核心策划圈")
    ws["A1"] = "核心策划圈（总策划 5/13 点名的高频沟通小圈子 + 葛九明，共 7 人）"
    ws["A1"].font = title_font
    ws.merge_cells("A1:C1")

    headers = ["姓名", "在分工表中的角色", "在核心圈的作用"]
    row = 3
    for i, h in enumerate(headers, 1):
        ws.cell(row=row, column=i, value=h)
    style_header(ws, row, len(headers))

    row += 1
    start = row
    for name, role, fn in CORE_CIRCLE:
        ws.cell(row=row, column=1, value=name)
        ws.cell(row=row, column=2, value=role)
        ws.cell(row=row, column=3, value=fn)
        row += 1
    end = row - 1
    style_body(ws, start, end, len(headers))
    for r in range(start, end+1):
        ws.row_dimensions[r].height = 30
    widths = [18, 22, 50]
    for i, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = w

    wb.save(path)
    print(f"[OK] Excel -> {path}")


# =============================================================================
# 生成 PPT
# =============================================================================
def set_cn_font(run, size=14, bold=False, color=None, name="微软雅黑"):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("eastAsia", "cs", "ascii", "hAnsi"):
        e = rPr.find(qn(f"a:{tag}"))
        if e is None:
            e = rPr.makeelement(qn(f"a:{tag}"), {"typeface": name})
            rPr.append(e)
        else:
            e.set("typeface", name)
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*hex_to_rgb(color))

def add_textbox(slide, left, top, width, height, text, size=14,
                bold=False, color="262626", align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_cn_font(run, size=size, bold=bold, color=color)
    return tb

def add_rect(slide, left, top, width, height, fill_hex,
             line_hex=None, text=None, text_size=14,
             text_bold=False, text_color="FFFFFF",
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             corner_radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(*hex_to_rgb(fill_hex))
    if line_hex is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = RGBColor(*hex_to_rgb(line_hex))
    if text is not None:
        tf = shp.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(36000)
        tf.margin_right = Emu(36000)
        tf.margin_top = Emu(18000)
        tf.margin_bottom = Emu(18000)
        tf.vertical_anchor = anchor
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        set_cn_font(run, size=text_size, bold=text_bold, color=text_color)
    return shp

def slide_title_bar(slide, title, subtitle=None, primary="1F4E79"):
    # Top color bar
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.85),
             primary)
    add_textbox(slide, Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.65),
                title, size=22, bold=True, color="FFFFFF",
                anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_textbox(slide, Inches(0.4), Inches(0.92), Inches(12.5),
                    Inches(0.35), subtitle, size=12, color="595959")

def build_ppt(path):
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ============ Slide 1 封面 ============
    s = prs.slides.add_slide(blank)
    # 背景
    add_rect(s, 0, 0, prs.slide_width, prs.slide_height, "0B2545")
    # 装饰条
    add_rect(s, Inches(0), Inches(2.7), Inches(13.333), Inches(0.08), "ED7D31")
    add_rect(s, Inches(0), Inches(5.0), Inches(13.333), Inches(0.04), "ED7D31")
    add_textbox(s, Inches(0.6), Inches(1.2), Inches(12), Inches(1.0),
                EVENT_TITLE, size=34, bold=True, color="FFFFFF")
    add_textbox(s, Inches(0.6), Inches(2.05), Inches(12), Inches(0.6),
                EVENT_SUB, size=18, color="F2F2F2")
    add_textbox(s, Inches(0.6), Inches(3.0), Inches(12), Inches(0.6),
                "志愿者职能分工方案 · 总策划工作底稿", size=22, bold=True, color="FFC000")
    add_textbox(s, Inches(0.6), Inches(3.7), Inches(12), Inches(0.5),
                f"时间：{EVENT_DATE}", size=14, color="D9D9D9")
    add_textbox(s, Inches(0.6), Inches(4.1), Inches(12), Inches(0.5),
                f"地点：{EVENT_VENUE}", size=14, color="D9D9D9")
    add_textbox(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
                "主办：北大经院上海校友会 · 复旦大学住房政策研究中心",
                size=12, color="BFBFBF")

    # ============ Slide 2 目录 ============
    s = prs.slides.add_slide(blank)
    slide_title_bar(s, "目录 / 总策划工作框架")
    items = [
        "01  整体目标与原则",
        "02  志愿者组织架构（向总策划负责）",
        "03  七大职能组职责（大项→小项）",
        "04  阶段任务矩阵（T-7 至 收尾撤场）",
        "05  议程时间轴分工",
        "06  应急预案（含技术/嘉宾/流程/安保）",
        "07  物资与对讲机频道",
        "08  关键检查节点与复盘",
    ]
    for i, it in enumerate(items):
        r = i // 2
        c = i % 2
        left = Inches(0.7 + c*6.2)
        top  = Inches(1.6 + r*1.1)
        add_rect(s, left, top, Inches(0.9), Inches(0.9), "1F4E79",
                 text=it.split()[0], text_size=20, text_bold=True,
                 corner_radius=True)
        add_textbox(s, left+Inches(1.1), top+Inches(0.15), Inches(5),
                    Inches(0.6), it[4:], size=16, color="262626",
                    anchor=MSO_ANCHOR.MIDDLE)

    # ============ Slide 3 目标与原则 ============
    s = prs.slides.add_slide(blank)
    slide_title_bar(s, "01  整体目标与工作原则")

    goals = [
        ("零失误执行", "500 人高净值圈层峰会，议程紧凑，任何环节零失误"),
        ("嘉宾体验",   "顶级嘉宾从抵达到离场全程被'看见、被服务'"),
        ("品牌呈现",   "直播+新媒体即时传播，会后 48h 形成传播声量"),
        ("应急可控",   "对任何可预见风险有预案、有责任人、有处置 SOP"),
    ]
    for i, (k, v) in enumerate(goals):
        top = Inches(1.5 + i*1.3)
        add_rect(s, Inches(0.6), top, Inches(2.6), Inches(1.0),
                 "C00000", text=k, text_size=18, text_bold=True)
        add_rect(s, Inches(3.3), top, Inches(9.4), Inches(1.0),
                 "F8F8F8", line_hex="D9D9D9", text=v, text_size=14,
                 text_color="262626", align=PP_ALIGN.LEFT)

    # ============ Slide 4 组织架构图 ============
    s = prs.slides.add_slide(blank)
    slide_title_bar(s, "02  志愿者组织架构", "总策划 → 执行总监(韦佳玉) → 7 大职能组")

    # 总策划
    add_rect(s, Inches(5.5), Inches(1.4), Inches(2.3), Inches(0.8),
             "0B2545", text="总策划 胡继刚（您）", text_size=16, text_bold=True,
             corner_radius=True)
    # 执行总监
    add_rect(s, Inches(5.5), Inches(2.5), Inches(2.3), Inches(0.7),
             "1F4E79", text="A 统筹协调组（韦佳玉 执行总监）", text_size=11,
             text_bold=True, corner_radius=True)
    # 连接线（用细矩形模拟）
    add_rect(s, Inches(6.6), Inches(2.2), Inches(0.05), Inches(0.3), "BFBFBF")
    add_rect(s, Inches(6.6), Inches(3.2), Inches(0.05), Inches(0.3), "BFBFBF")
    add_rect(s, Inches(0.7), Inches(3.5), Inches(11.95), Inches(0.04), "BFBFBF")

    # 6 个子组
    sub_groups = GROUPS[1:]  # B-G
    n = len(sub_groups)
    total_w = 12.0
    box_w = 1.85
    gap = (total_w - box_w*n) / (n-1)
    for i, g in enumerate(sub_groups):
        left = Inches(0.7 + i*(box_w+gap))
        # 引线
        add_rect(s, left + Inches(box_w/2 - 0.025), Inches(3.5),
                 Inches(0.05), Inches(0.4), "BFBFBF")
        add_rect(s, left, Inches(3.9), Inches(box_w), Inches(0.6),
                 g["color"], text=f"{g['code']} {g['name']}",
                 text_size=13, text_bold=True, corner_radius=True)
        # 人数 + 负责人
        add_rect(s, left, Inches(4.55), Inches(box_w), Inches(1.6),
                 "F8F8F8", line_hex=g["color"])
        add_textbox(s, left+Inches(0.1), Inches(4.6), Inches(box_w-0.2),
                    Inches(0.4), f"负责人：{g['leader']}",
                    size=10, color="262626", bold=True)
        add_textbox(s, left+Inches(0.1), Inches(4.95), Inches(box_w-0.2),
                    Inches(0.3), f"人数：{g['size']}", size=10, color="595959")
        add_textbox(s, left+Inches(0.1), Inches(5.25), Inches(box_w-0.2),
                    Inches(0.9), g["duty_summary"], size=9, color="404040")

    add_textbox(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.5),
                "原则：每组设组长 1 人 → 直接对接执行总监 A 组 → A 组对您负责；对讲机分 3 频道（指挥/会务/后勤）",
                size=12, color="C00000", bold=True)

    # ============ Slide 5-11 每组职责详情 ============
    for g in GROUPS:
        s = prs.slides.add_slide(blank)
        slide_title_bar(s, f"03  {g['code']} {g['name']} · 职责细分（大项→小项）",
                        f"负责人：{g['leader']}  |  建议人数：{g['size']}",
                        primary=g["color"])

        add_rect(s, Inches(0.6), Inches(1.5), Inches(12.2), Inches(0.7),
                 "F2F2F2", line_hex=g["color"],
                 text=f"主要职责：{g['duty_summary']}",
                 text_size=14, text_color="262626", align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.MIDDLE)

        n = len(g["sub_duties"])
        row_h = (6.4 - 2.3) / max(n, 1)
        for i, (sub_name, sub_desc) in enumerate(g["sub_duties"]):
            top = Inches(2.4 + i*row_h)
            add_rect(s, Inches(0.6), top, Inches(2.6), Inches(row_h - 0.1),
                     g["color"], text=f"{i+1}. {sub_name}",
                     text_size=13, text_bold=True)
            add_rect(s, Inches(3.3), top, Inches(9.5), Inches(row_h - 0.1),
                     "FFFFFF", line_hex="D9D9D9", text=sub_desc,
                     text_size=12, text_color="262626", align=PP_ALIGN.LEFT)

        add_textbox(s, Inches(0.6), Inches(7.0), Inches(12), Inches(0.4),
                    f"对接：A 统筹协调组   |   对讲机频道：见物资页",
                    size=10, color="808080")

    # ============ 本轮人员调整说明（v4） ============
    s = prs.slides.add_slide(blank)
    slide_title_bar(s, "本轮人员调整说明（总策划 5/14 v4 调整）",
                    "原则：应急组就位/技术组组长换律师/茶歇晚宴最适配人选/A 组职责再细分")

    # 左侧：调整动作
    add_rect(s, Inches(0.55), Inches(1.4), Inches(6.0), Inches(0.5),
             "C00000", text="本轮 8 项调整", text_size=14, text_bold=True)
    moves = [
        ("组长", "刘严：F 副 → G 组长（应急/机动调配）"),
        ("组长", "蔡萍：G 急救主岗 → G 副组长（兼急救★）"),
        ("组长", "张卓：G 副 → E 组长（律师强协调，统筹技术）"),
        ("人员", "Winny：B → F 副组长（茶歇/晚宴主管 ★最适合）"),
        ("人员", "吕志翔：E → C（投融资非技术，回归会务）"),
        ("人员", "马磊：D → G（机动+媒体应急口径）"),
        ("职责", "李振春：文件秘书 → 科创赞助专项（聚焦）"),
        ("职责", "韦佳玉：执行总监 + 文件总管（接管李振春原职责）"),
    ]
    row_h = 0.6
    for i, (k, v) in enumerate(moves):
        y = Inches(2.0 + i*row_h)
        color = {"组长": "C00000", "人员": "ED7D31", "职责": "1F4E79"}[k]
        add_rect(s, Inches(0.55), y, Inches(0.9), Inches(row_h-0.08),
                 color, text=k, text_size=10, text_bold=True)
        add_rect(s, Inches(1.5), y, Inches(5.05), Inches(row_h-0.08),
                 "FFFFFF", line_hex="D9D9D9", text=v,
                 text_size=10, text_color="262626", align=PP_ALIGN.LEFT)

    # 右侧：人数变化
    add_rect(s, Inches(6.85), Inches(1.4), Inches(6.0), Inches(0.5),
             "1F4E79", text="人数变化（27 志愿者 + 总策划 = 28 人）",
             text_size=14, text_bold=True)
    headers = ["组", "v3", "v4", "变化", "组长 / 副组长"]
    widths  = [1.2, 0.7, 0.7, 0.9, 2.5]
    top = Inches(2.0)
    head_h = 0.42
    x = 6.85
    for h, w in zip(headers, widths):
        add_rect(s, Inches(x), top, Inches(w), Inches(head_h),
                 "2E75B6", text=h, text_size=10, text_bold=True)
        x += w
    delta_data = [
        ("A 统筹",  "3", "3", "—",   "胡继刚 / 韦佳玉"),
        ("B 接待",  "6", "5", "↓ 1", "王胜 / 朱铭喆"),
        ("C 会务",  "5", "6", "↑ 1", "陈潇 Kelly / 王珏"),
        ("D 媒体",  "5", "4", "↓ 1", "葛九明 / 张蒙"),
        ("E 技术",  "3", "3", "—",   "★ 张卓 / 随圣博"),
        ("F 后勤",  "3", "3", "—",   "朱俊峰 / ★ Winny"),
        ("G 应急",  "3", "4", "↑ 1", "★ 刘严 / ★ 蔡萍"),
    ]
    row_h = 0.55
    for ri, row in enumerate(delta_data):
        y = Inches(2.0 + head_h + ri*row_h)
        x = 6.85
        # 颜色提示
        for w, v, idx in zip(widths, row, range(5)):
            fill = "FFFFFF" if ri % 2 == 0 else "F2F8FD"
            tcolor = "262626"
            tbold = False
            if idx == 3 and "↓" in v:
                fill = "FFE4E4"; tcolor = "C00000"; tbold = True
            elif idx == 3 and "↑" in v:
                fill = "E4FFE4"; tcolor = "548235"; tbold = True
            add_rect(s, Inches(x), y, Inches(w), Inches(row_h),
                     fill, line_hex="D9D9D9", text=v,
                     text_size=10, text_color=tcolor, text_bold=tbold,
                     align=PP_ALIGN.CENTER if idx<4 else PP_ALIGN.LEFT,
                     anchor=MSO_ANCHOR.MIDDLE)
            x += w

    # ============ 实际人员分配总表 ============
    s = prs.slides.add_slide(blank)
    total_n = len(ASSIGNMENT)
    slide_title_bar(s, f"实际人员分配（基于总策划最新调整）",
                    f"共 {total_n} 人 / 7 组；本轮调整：+葛九明、蒋珊→D、吕志翔/高辰辰→E、彭常丽→F、王胜→B组长、陈潇→C组长")

    # 按组分块布局
    blocks = {}
    for code, post, name, pref, bg, reason in ASSIGNMENT:
        blocks.setdefault(code, []).append((post, name, pref))

    color_map = {g["code"]: g["color"] for g in GROUPS}
    name_map  = {g["code"]: g["name"] for g in GROUPS}

    # 4 列 x 2 行布局（更均衡）
    # 行1: A B C D / 行2: E F G
    grid_top1, grid_top2 = 1.45, 4.45
    grid_h = 2.85
    cols_top = ["A", "B", "C", "D"]
    cols_bot = ["E", "F", "G"]
    n_top, n_bot = len(cols_top), len(cols_bot)
    margin = 0.4
    avail_w = 13.333 - 2*margin
    col_w_top = (avail_w - (n_top-1)*0.15) / n_top
    col_w_bot = (avail_w - (n_bot-1)*0.15) / n_bot

    def draw_block(x, y, w, h, code):
        add_rect(s, Inches(x), Inches(y), Inches(w), Inches(0.42),
                 color_map[code],
                 text=f"{code} {name_map[code]}（{len(blocks[code])}人）",
                 text_size=12, text_bold=True)
        add_rect(s, Inches(x), Inches(y+0.42), Inches(w), Inches(h-0.42),
                 "FFFFFF", line_hex=color_map[code])
        members = blocks[code]
        line_h = (h - 0.55) / max(len(members), 1)
        for i, (post, name, pref) in enumerate(members):
            top = Inches(y + 0.5 + i*line_h)
            # 标记组长/副组长
            mark = ""
            if "组长" in post and "副" not in post:
                mark = "★"
            elif "副组长" in post:
                mark = "☆"
            elif "总策划" in post or "执行总监" in post:
                mark = "★"
            add_textbox(s, Inches(x+0.1), top, Inches(w-0.2),
                        Inches(line_h),
                        f"{mark}  {name}  -  {post.split('/')[0].strip()}",
                        size=9, color="262626", anchor=MSO_ANCHOR.MIDDLE)

    for i, code in enumerate(cols_top):
        x = margin + i*(col_w_top + 0.15)
        draw_block(x, grid_top1, col_w_top, grid_h, code)
    for i, code in enumerate(cols_bot):
        x = margin + i*(col_w_bot + 0.15)
        draw_block(x, grid_top2, col_w_bot, grid_h, code)

    # ============ 5/13 优化建议筛选（采纳/调整/暂缓） ============
    s = prs.slides.add_slide(blank)
    slide_title_bar(s, "5/13 优化建议筛选（共 19 条 → 采纳 16 / 调整 1 / 暂缓 2）",
                    "处置原则：直击去年痛点的全部采纳；专业度不足或越界的暂缓",
                    primary="C00000")

    # 表头
    headers_w = [(1.4, "建议人"), (5.0, "建议内容"), (1.3, "处置"),
                 (1.0, "组"), (4.4, "落地动作")]
    top = Inches(1.5)
    head_h = 0.42
    x = 0.55
    for w, h in headers_w:
        add_rect(s, Inches(x), top, Inches(w), Inches(head_h),
                 "1F4E79", text=h, text_size=10, text_bold=True)
        x += w

    status_color = {
        "采纳": "548235",
        "调整后采纳": "ED7D31",
        "暂缓": "808080",
        "暂缓→单列待办": "808080",
    }
    rows_per_slide = len(OPT_SUGGESTIONS)
    row_h = (7.2 - 1.5 - 0.42) / rows_per_slide
    for ri, (sug, content, status, resp, act) in enumerate(OPT_SUGGESTIONS):
        y = Inches(1.5 + head_h + ri*row_h)
        x = 0.55
        cells = [(headers_w[0][0], sug, "262626", "FFFFFF" if ri%2==0 else "F8F8F8"),
                 (headers_w[1][0], content, "262626", "FFFFFF" if ri%2==0 else "F8F8F8"),
                 (headers_w[2][0], status, "FFFFFF", status_color.get(status, "808080")),
                 (headers_w[3][0], resp, "262626", "FFFFFF" if ri%2==0 else "F8F8F8"),
                 (headers_w[4][0], act, "262626", "FFFFFF" if ri%2==0 else "F8F8F8")]
        for w, val, tcolor, fill in cells:
            bold = (val == status and val in status_color)
            align = PP_ALIGN.CENTER if val == status else PP_ALIGN.LEFT
            add_rect(s, Inches(x), y, Inches(w), Inches(row_h),
                     fill, line_hex="D9D9D9", text=val,
                     text_size=8, text_bold=bold, text_color=tcolor,
                     align=align, anchor=MSO_ANCHOR.MIDDLE)
            x += w

    # ============ 核心策划圈 + 组长 A/B 角 ============
    s = prs.slides.add_slide(blank)
    slide_title_bar(s, "核心策划圈 + 组长 A/B 角备份制度",
                    "采纳王胜建议：组长备份；总策划 5/13 点名核心搭子")

    # 左：核心策划圈
    add_rect(s, Inches(0.55), Inches(1.4), Inches(6.0), Inches(0.5),
             "0B2545", text="核心策划圈（7 人 · 总策划直管高频沟通）",
             text_size=13, text_bold=True)
    for i, (name, role, fn) in enumerate(CORE_CIRCLE):
        top = Inches(2.0 + i*0.75)
        add_rect(s, Inches(0.55), top, Inches(1.7), Inches(0.65),
                 "2E75B6", text=name, text_size=11, text_bold=True)
        add_rect(s, Inches(2.3), top, Inches(1.6), Inches(0.65),
                 "F2F2F2", line_hex="2E75B6", text=role,
                 text_size=10, text_color="262626")
        add_rect(s, Inches(3.95), top, Inches(2.6), Inches(0.65),
                 "FFFFFF", line_hex="2E75B6", text=fn,
                 text_size=9, text_color="404040", align=PP_ALIGN.LEFT)

    # 右：组长 A/B 角
    add_rect(s, Inches(6.85), Inches(1.4), Inches(6.0), Inches(0.5),
             "C00000", text="组长 A/B 角备份（采纳王胜建议）",
             text_size=13, text_bold=True)
    headers = ["组", "A 角组长", "B 角备份"]
    widths  = [1.2, 1.8, 1.8]
    top = Inches(2.0)
    head_h = 0.42
    x = 6.85
    for h, w in zip(headers, widths):
        add_rect(s, Inches(x), top, Inches(w), Inches(head_h),
                 "ED7D31", text=h, text_size=11, text_bold=True)
        x += w
    # remaining width for note
    add_rect(s, Inches(x), top, Inches(6.0-sum(widths)), Inches(head_h),
             "ED7D31", text="备注", text_size=11, text_bold=True)

    row_h = 0.55
    for ri, (grp, a, b, note) in enumerate(BACKUP_ROLES):
        y = Inches(2.0 + head_h + ri*row_h)
        x = 6.85
        for w, v in zip(widths + [6.0-sum(widths)], [grp, a, b, note]):
            fill = "FFFFFF" if ri % 2 == 0 else "FFF7EC"
            add_rect(s, Inches(x), y, Inches(w), Inches(row_h),
                     fill, line_hex="D9D9D9", text=v,
                     text_size=9, text_color="262626",
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
            x += w

    # ============ Milestone 时间表 ============
    s = prs.slides.add_slide(blank)
    slide_title_bar(s, "Milestone 时间表（节点提前 1-2 天 / 韦佳玉建议）",
                    f"共 {len(MILESTONES)} 个里程碑 / T-9 ~ T+1")

    headers = ["时间节点", "关键里程碑", "责任组"]
    widths  = [1.7, 9.4, 1.6]
    top = Inches(1.5)
    head_h = 0.42
    x = 0.55
    for h, w in zip(headers, widths):
        add_rect(s, Inches(x), top, Inches(w), Inches(head_h),
                 "1F4E79", text=h, text_size=11, text_bold=True)
        x += w

    row_h = (7.2 - 1.5 - 0.42) / len(MILESTONES)
    for ri, (t, ms, resp) in enumerate(MILESTONES):
        y = Inches(1.5 + head_h + ri*row_h)
        x = 0.55
        # T-0/T-1/T-7 用不同色突出
        if t.startswith("T-0"):
            fill = "FFE4E4"
        elif t.startswith("T-1"):
            fill = "FFF0E4"
        elif t.startswith("T+1"):
            fill = "E4F0FF"
        else:
            fill = "FFFFFF" if ri % 2 == 0 else "F8F8F8"
        for w, v in zip(widths, [t, ms, resp]):
            add_rect(s, Inches(x), y, Inches(w), Inches(row_h),
                     fill, line_hex="D9D9D9", text=v,
                     text_size=8, text_color="262626",
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
            x += w

    # ============ 嘉宾发言时间控制 SOP ============
    s = prs.slides.add_slide(blank)
    slide_title_bar(s, "嘉宾发言时间控制 SOP（去年最大事故）",
                    "采纳张卓建议：从第 1 位讲者就开始的超时绝不能重演",
                    primary="C00000")

    headers = ["时间点", "动作", "责任组"]
    widths  = [2.4, 8.7, 1.6]
    top = Inches(1.5)
    head_h = 0.5
    x = 0.55
    for h, w in zip(headers, widths):
        add_rect(s, Inches(x), top, Inches(w), Inches(head_h),
                 "C00000", text=h, text_size=12, text_bold=True)
        x += w

    row_h = (7.2 - 1.5 - 0.5) / len(TIME_SOP)
    for ri, (t, act, resp) in enumerate(TIME_SOP):
        y = Inches(1.5 + head_h + ri*row_h)
        x = 0.55
        # 现场动作高亮
        is_live = t.startswith("T-0")
        fill = "FFF7F7" if is_live else ("FFFFFF" if ri%2==0 else "F8F8F8")
        for w, v in zip(widths, [t, act, resp]):
            add_rect(s, Inches(x), y, Inches(w), Inches(row_h),
                     fill, line_hex="D9D9D9", text=v,
                     text_size=10, text_color="262626",
                     text_bold=is_live,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
            x += w

    # ============ 协助人员对接清单 ============
    s = prs.slides.add_slide(blank)
    slide_title_bar(s, "协助人员对接清单（全员人手一份）",
                    "采纳张卓建议：现场遇问题先找A角→找B角→才找韦佳玉/胡继刚")

    headers = ["问题类别 / 你遇到的事情", "A 角（先找）", "B 角（找不到时）"]
    widths  = [5.4, 3.5, 3.8]
    top = Inches(1.5)
    head_h = 0.45
    x = 0.55
    for h, w in zip(headers, widths):
        add_rect(s, Inches(x), top, Inches(w), Inches(head_h),
                 "548235", text=h, text_size=11, text_bold=True)
        x += w

    row_h = (7.2 - 1.5 - 0.45) / len(LIAISON_LIST)
    for ri, (cat, who, backup) in enumerate(LIAISON_LIST):
        y = Inches(1.5 + head_h + ri*row_h)
        x = 0.55
        # 重要分类高亮
        is_critical = ("医疗" in cat or "重大" in cat or "舆情" in cat)
        fill = "FFE4E4" if is_critical else ("FFFFFF" if ri%2==0 else "F4F8EE")
        for w, v in zip(widths, [cat, who, backup]):
            add_rect(s, Inches(x), y, Inches(w), Inches(row_h),
                     fill, line_hex="D9D9D9", text=v,
                     text_size=9, text_color="262626",
                     text_bold=is_critical,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
            x += w

    # ============ 5/13 会议结论 & 待办 ============
    s = prs.slides.add_slide(blank)
    slide_title_bar(s, "5/13 志愿者讨论会  会议结论 & 待办")

    # 左：会议结论
    add_rect(s, Inches(0.55), Inches(1.4), Inches(6.0), Inches(0.5),
             "1F4E79", text="会议结论（4 条优化共识）",
             text_size=14, text_bold=True)
    conclusions = [
        ("分工与协作", "各组按任务量配人，组长负责内部协调；设机动应急G组跨组顶人；各组建独立沟通群"),
        ("流程效率",   "任务量与 Milestone 提前细化；所有节点提前 1-2 天预留缓冲"),
        ("交响乐环节", "扩大宣传，发往届视频预热，与嘉宾/观众多互动"),
        ("协作机制",   "活动前组织一次组长见面会，互相熟悉、快速找到对接人"),
    ]
    for i, (k, v) in enumerate(conclusions):
        top = Inches(2.0 + i*1.25)
        add_rect(s, Inches(0.55), top, Inches(1.5), Inches(1.15),
                 "2E75B6", text=k, text_size=11, text_bold=True)
        add_rect(s, Inches(2.1), top, Inches(4.45), Inches(1.15),
                 "F2F8FD", line_hex="2E75B6", text=v,
                 text_size=10, text_color="262626",
                 align=PP_ALIGN.LEFT)

    # 右：待办
    add_rect(s, Inches(6.85), Inches(1.4), Inches(6.0), Inches(0.5),
             "C00000", text="待办（总策划 + 各组组长）",
             text_size=14, text_bold=True)
    todos = [
        ("胡继刚",  "形成详细分工表（本表）；私信招募 ①近期有时间者 ②有会务经验者 ③有赞助经验/资源者"),
        ("各组长","内部完成细化分工；建立独立沟通群；T-7 前提交本组 Milestone"),
        ("会议",   "另开 2 次会：① 各组任务进一步细化  ② 线上彩排"),
        ("特别行动","组长见面会（A+B/C/D/E/F/G 组长 共 8 人）建议 T-10 前完成"),
        ("交响乐", "F 组对接陆家嘴乐团 + D 组准备往届视频预热"),
    ]
    for i, (k, v) in enumerate(todos):
        top = Inches(2.0 + i*1.0)
        add_rect(s, Inches(6.85), top, Inches(1.5), Inches(0.9),
                 "ED7D31", text=k, text_size=11, text_bold=True)
        add_rect(s, Inches(8.4), top, Inches(4.45), Inches(0.9),
                 "FFF7EC", line_hex="ED7D31", text=v,
                 text_size=10, text_color="262626",
                 align=PP_ALIGN.LEFT)

    # ============ Slide 阶段任务矩阵（拆分两张） ============
    # 第一张：T-7 / T-1 / 当日布置
    def add_phase_matrix(slide, phase_indices, title):
        slide_title_bar(slide, title, "横向：组别  |  纵向：阶段")
        n_phase = len(phase_indices)
        n_group = len(GROUPS)
        total_left = 0.5
        total_top  = 1.5
        total_w    = 12.3
        total_h    = 5.7
        # 表头：组别行
        head_h = 0.55
        first_col_w = 1.6
        col_w = (total_w - first_col_w) / n_group
        row_h = (total_h - head_h) / n_phase

        # 空角
        add_rect(slide, Inches(total_left), Inches(total_top),
                 Inches(first_col_w), Inches(head_h),
                 "1F4E79", text="阶段 \\ 组别", text_size=11, text_bold=True)
        # 组别表头
        for gi, g in enumerate(GROUPS):
            add_rect(slide, Inches(total_left + first_col_w + gi*col_w),
                     Inches(total_top), Inches(col_w), Inches(head_h),
                     g["color"], text=f"{g['code']} {g['name']}",
                     text_size=10, text_bold=True)
        # 行
        for ri, pi in enumerate(phase_indices):
            when, name = PHASES[pi]
            top = Inches(total_top + head_h + ri*row_h)
            add_rect(slide, Inches(total_left), top, Inches(first_col_w),
                     Inches(row_h), "F2F2F2", line_hex="BFBFBF",
                     text=f"{name}\n{when}", text_size=10, text_bold=True,
                     text_color="1F4E79")
            for gi, g in enumerate(GROUPS):
                task = PHASE_MATRIX.get((pi, g["code"]), "")
                add_rect(slide, Inches(total_left + first_col_w + gi*col_w),
                         top, Inches(col_w), Inches(row_h),
                         "FFFFFF", line_hex="D9D9D9", text=task,
                         text_size=8, text_color="262626",
                         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

    s = prs.slides.add_slide(blank)
    add_phase_matrix(s, [0,1,2], "04  阶段任务矩阵（一）会前筹备")

    s = prs.slides.add_slide(blank)
    add_phase_matrix(s, [3,4,5], "04  阶段任务矩阵（二）当日执行 + 收尾")

    # ============ 议程时间轴分工 ============
    s = prs.slides.add_slide(blank)
    slide_title_bar(s, "05  议程时间轴 × 责任组 × 关键动作")

    headers = ["时间", "议程节点", "主责组", "关键动作 / 物料"]
    widths = [1.4, 4.4, 1.4, 5.5]
    top = Inches(1.5)
    head_h = 0.5
    left_start = 0.55
    # header
    x = left_start
    for h, w in zip(headers, widths):
        add_rect(s, Inches(x), top, Inches(w), Inches(head_h),
                 "1F4E79", text=h, text_size=12, text_bold=True)
        x += w
    # rows
    row_h = (7.3 - 1.5 - 0.5) / len(AGENDA)
    for ri, (t, name, gr, action) in enumerate(AGENDA):
        y = Inches(1.5 + head_h + ri*row_h)
        x = left_start
        values = [t, name, gr, action]
        for v, w in zip(values, widths):
            fill = "FFFFFF" if ri % 2 == 0 else "F8F8F8"
            add_rect(s, Inches(x), y, Inches(w), Inches(row_h),
                     fill, line_hex="D9D9D9", text=v, text_size=9,
                     text_color="262626", align=PP_ALIGN.LEFT,
                     anchor=MSO_ANCHOR.MIDDLE)
            x += w

    # ============ 应急预案（分两张） ============
    def add_contingency_slide(slide, items, title):
        slide_title_bar(slide, title)
        headers = ["类别","风险","概率","影响","责任组","处置预案"]
        widths  = [1.2, 2.6, 0.7, 0.7, 1.3, 6.2]
        top = Inches(1.5)
        head_h = 0.5
        x = 0.55
        for h, w in zip(headers, widths):
            add_rect(slide, Inches(x), top, Inches(w), Inches(head_h),
                     "C00000", text=h, text_size=11, text_bold=True)
            x += w
        row_h = (7.2 - 1.5 - 0.5) / max(len(items), 1)
        for ri, item in enumerate(items):
            y = Inches(1.5 + head_h + ri*row_h)
            x = 0.55
            for v, w in zip(item, widths):
                fill = "FFFFFF" if ri % 2 == 0 else "FFF7F7"
                add_rect(slide, Inches(x), y, Inches(w), Inches(row_h),
                         fill, line_hex="D9D9D9", text=v, text_size=8,
                         text_color="262626", align=PP_ALIGN.LEFT,
                         anchor=MSO_ANCHOR.MIDDLE)
                x += w

    half = (len(CONTINGENCIES)+1)//2
    s = prs.slides.add_slide(blank)
    add_contingency_slide(s, CONTINGENCIES[:half],
                          "06  应急预案（一）嘉宾 / 技术 / 流程")
    s = prs.slides.add_slide(blank)
    add_contingency_slide(s, CONTINGENCIES[half:],
                          "06  应急预案（二）安全 / 舆情 / 后勤 / 晚宴")

    # ============ 物资清单 + 对讲机频道 ============
    s = prs.slides.add_slide(blank)
    slide_title_bar(s, "07  物资清单 与 对讲机频道")

    # 物资表
    headers = ["类别","物品","数量","负责组","到位时间"]
    widths  = [0.9, 2.5, 2.2, 1.1, 1.6]
    top = Inches(1.4)
    head_h = 0.45
    x = 0.55
    for h, w in zip(headers, widths):
        add_rect(s, Inches(x), top, Inches(w), Inches(head_h),
                 "548235", text=h, text_size=11, text_bold=True)
        x += w
    row_h = 0.35
    for ri, item in enumerate(MATERIALS):
        y = Inches(1.4 + head_h + ri*row_h)
        x = 0.55
        for v, w in zip(item, widths):
            fill = "FFFFFF" if ri % 2 == 0 else "F4F8EE"
            add_rect(s, Inches(x), y, Inches(w), Inches(row_h),
                     fill, line_hex="D9D9D9", text=v, text_size=9,
                     text_color="262626", align=PP_ALIGN.LEFT,
                     anchor=MSO_ANCHOR.MIDDLE)
            x += w

    # 对讲机频道
    add_rect(s, Inches(9.5), Inches(1.4), Inches(3.4), Inches(0.45),
             "1F4E79", text="对讲机频道", text_size=12, text_bold=True)
    channels = [
        ("CH1 指挥频道",  "总策划 / A 组 / 各组组长"),
        ("CH2 会务频道",  "C 会务 + E 技术（流程切换）"),
        ("CH3 后勤频道",  "B 接待 + F 后勤 + G 安全"),
        ("公共哨语",      "'1分钟到点'、'换场'、'医疗'、'静默'"),
    ]
    for i, (k, v) in enumerate(channels):
        y = Inches(1.85 + i*0.7)
        add_rect(s, Inches(9.5), y, Inches(1.3), Inches(0.65),
                 "2E75B6", text=k, text_size=10, text_bold=True)
        add_rect(s, Inches(10.8), y, Inches(2.1), Inches(0.65),
                 "FFFFFF", line_hex="2E75B6", text=v, text_size=10,
                 text_color="262626", align=PP_ALIGN.LEFT)

    # ============ 关键检查节点 ============
    s = prs.slides.add_slide(blank)
    slide_title_bar(s, "08  关键检查节点（Checkpoint）")

    checkpoints = [
        ("T-3 日 18:00",  "所有嘉宾 PPT/演讲稿到位；议程脚本 V2 锁定；通稿初稿"),
        ("T-1 日 14:00",  "现场走台彩排；全员到岗培训；备用设备压力测试"),
        ("当日 09:00",    "现场指挥部成立；对讲机分发；安保/医疗到位"),
        ("当日 12:00",    "签到台开台；VIP车队首班；直播试推流"),
        ("当日 13:25",    "全员就位静默 5 分钟；总策划广播开场倒计时"),
        ("当日 15:40",    "音乐会前节点 - 茶歇/伴手礼/媒体合影同步启动"),
        ("当日 18:10",    "颁奖准备 - 奖杯/获奖机构到台口；晚宴席位最终核对"),
        ("当日 18:30",    "晚宴开席 - 主会场清场启动撤场"),
        ("次日 10:00",    "通稿发布；志愿者复盘会；问题与改进归档"),
    ]
    for i, (t, v) in enumerate(checkpoints):
        r, c = i // 3, i % 3
        left = Inches(0.55 + c*4.3)
        top  = Inches(1.6 + r*1.9)
        add_rect(s, left, top, Inches(4.05), Inches(0.55),
                 "ED7D31", text=t, text_size=12, text_bold=True,
                 corner_radius=True)
        add_rect(s, left, top + Inches(0.55), Inches(4.05), Inches(1.2),
                 "FFF7EC", line_hex="ED7D31", text=v, text_size=11,
                 text_color="262626", align=PP_ALIGN.LEFT)

    # ============ 结尾 ============
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, prs.slide_width, prs.slide_height, "0B2545")
    add_textbox(s, Inches(0.6), Inches(2.4), Inches(12), Inches(1.2),
                "万事俱备，只待 5·22。", size=44, bold=True,
                color="FFFFFF", align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.6), Inches(3.8), Inches(12), Inches(0.6),
                "所有志愿者的努力，都是为了让 500 位嘉宾在一滴水里看见 AI 的未来。",
                size=18, color="FFC000", align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.6), Inches(5.0), Inches(12), Inches(0.5),
                "—— 总策划",
                size=14, color="D9D9D9", align=PP_ALIGN.CENTER)

    prs.save(path)
    print(f"[OK] PPT -> {path}")


if __name__ == "__main__":
    build_excel("/workspace/2026峰会_志愿者分工表.xlsx")
    build_ppt("/workspace/2026峰会_志愿者分工方案.pptx")
