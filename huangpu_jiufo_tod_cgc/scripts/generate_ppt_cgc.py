# -*- coding: utf-8 -*-
"""生成《广州黄埔区九佛TOD全球自贸365街区策划方案》优化商务风 PPT(16:9)。"""
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Cm, Emu, Pt

import content_cgc as C

OUT_DIR = os.path.join(os.path.dirname(__file__), "..", "deliverables")
FONT = "微软雅黑"

NAVY = RGBColor(0x0F, 0x2A, 0x4A)       # 主色 深藏青
BLUE = RGBColor(0x1F, 0x4E, 0x79)       # 次色 蓝
GOLD = RGBColor(0xC7, 0x9A, 0x3B)       # 强调 金
LIGHT = RGBColor(0xEE, 0xF2, 0xF8)      # 浅底
LIGHTGOLD = RGBColor(0xF6, 0xEE, 0xDA)
GRAY = RGBColor(0x5A, 0x5A, 0x5A)
DARK = RGBColor(0x22, 0x2B, 0x35)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SW, SH = Cm(33.867), Cm(19.05)


def _set_ea(run, font):
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:latin"):
        existing = rPr.find(qn(tag))
        if existing is None:
            existing = rPr.makeelement(qn(tag), {})
            rPr.append(existing)
        existing.set("typeface", font)


def _run(p, text, size, bold=False, color=DARK, font=FONT):
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    _set_ea(r, font)
    return r


def _tf(shape):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    return tf


def textbox(slide, l, t, w, h, anchor=None):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = _tf(box)
    if anchor:
        tf.vertical_anchor = anchor
    return tf


def rect(slide, l, t, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, l, t, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def content_header(slide, section_no, section_name, title):
    """标准内容页页眉:左侧色条+编号,标题,下方金色分隔线。"""
    rect(slide, Cm(0), Cm(0), Cm(0.35), SH, GOLD)
    tf = textbox(slide, Cm(1.4), Cm(0.75), Cm(6), Cm(1.2))
    _run(tf.paragraphs[0], f"0{section_no}  {section_name}", 12, bold=True, color=GOLD)
    tf = textbox(slide, Cm(1.4), Cm(1.55), Cm(31), Cm(1.6))
    _run(tf.paragraphs[0], title, 24, bold=True, color=NAVY)
    rect(slide, Cm(1.5), Cm(3.15), Cm(6.5), Cm(0.09), GOLD)


def add_bullets(slide, items, left, top, width, height, size=15,
                gap=10, color=DARK, bullet_color=GOLD):
    tf = textbox(slide, left, top, width, height)
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        _run(p, "▍ ", size, bold=True, color=bullet_color)
        _run(p, it, size, color=color)
        p.space_after = Pt(gap)
        p.line_spacing = 1.15
    return tf


def card_grid(slide, cards, top, cols=2, left=Cm(1.4), total_w=Cm(31.0),
              row_h=Cm(3.0), h_gap=Cm(0.5), v_gap=Cm(0.5),
              title_size=15, body_size=12, accent=GOLD, tinted=True):
    """cards: list of (title, body) 或 (title, subtitle, body)。"""
    n = len(cards)
    rows = (n + cols - 1) // cols
    card_w = Emu(int((total_w - h_gap * (cols - 1)) / cols))
    for i, card in enumerate(cards):
        r, c = divmod(i, cols)
        l = Emu(int(left + c * (card_w + h_gap)))
        t = Emu(int(top + r * (row_h + v_gap)))
        bg = LIGHT if tinted else WHITE
        box = rect(slide, l, t, card_w, row_h, bg)
        rect(slide, l, t, Cm(0.16), row_h, accent)
        tf = _tf(box)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = Pt(12)
        tf.margin_right = Pt(10)
        tf.margin_top = Pt(9)
        if len(card) == 3:
            title, sub, body = card
            p = tf.paragraphs[0]
            _run(p, title, title_size + 3, bold=True, color=accent)
            _run(p, "  " + sub, title_size - 3, bold=True, color=NAVY)
        else:
            title, body = card
            p = tf.paragraphs[0]
            _run(p, title, title_size, bold=True, color=NAVY)
        p.space_after = Pt(5)
        pb = tf.add_paragraph()
        _run(pb, body, body_size, color=GRAY)
        pb.line_spacing = 1.12
    return rows


def make_table(slide, headers, rows, left, top, width, col_widths,
               header_size=11, body_size=9.5, row_h=Cm(0.9), header_h=Cm(0.95),
               zebra=True):
    nrows = len(rows) + 1
    ncols = len(headers)
    table = slide.shapes.add_table(nrows, ncols, left, top, width,
                                   Emu(int(header_h + row_h * len(rows)))).table
    table.first_row = False
    table.horz_banding = False
    total = sum(col_widths)
    for j, cw in enumerate(col_widths):
        table.columns[j].width = Emu(int(width * cw / total))
    table.rows[0].height = Emu(int(header_h))
    for i in range(1, nrows):
        table.rows[i].height = Emu(int(row_h))
    # header
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Pt(5); cell.margin_right = Pt(5)
        cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _run(p, h, header_size, bold=True, color=WHITE)
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if (zebra and i % 2 == 0) else WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Pt(5); cell.margin_right = Pt(5)
            cell.margin_top = Pt(1); cell.margin_bottom = Pt(1)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j == 0 else PP_ALIGN.LEFT
            _run(p, str(val), body_size, color=DARK)
    return table


# ============================================================= 幻灯片
def cover(prs):
    s = blank(prs)
    rect(s, Cm(0), Cm(0), SW, SH, NAVY)
    rect(s, Cm(0), Cm(11.2), SW, Cm(0.12), GOLD)
    tf = textbox(s, Cm(2.2), Cm(4.4), Cm(29.5), Cm(4.5))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, C.PROJECT_NAME, 30, bold=True, color=WHITE)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(14)
    _run(p2, C.SLOGAN, 20, bold=True, color=GOLD)
    tf = textbox(s, Cm(2.2), Cm(11.8), Cm(29.5), Cm(3))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, C.SUB_SLOGAN, 14, color=LIGHT)
    tf = textbox(s, Cm(2.2), Cm(15.6), Cm(29.5), Cm(2.5))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, f"{C.OWNER_ORG}   {C.OWNER_PERSON}", 15, bold=True, color=WHITE)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(6)
    _run(p2, C.DOC_DATE, 13, color=LIGHT)


def contents(prs):
    s = blank(prs)
    rect(s, Cm(0), Cm(0), Cm(11.5), SH, NAVY)
    tf = textbox(s, Cm(1.4), Cm(6.2), Cm(9.5), Cm(4))
    _run(tf.paragraphs[0], "CONTENTS", 30, bold=True, color=GOLD)
    p = tf.add_paragraph(); p.space_before = Pt(6)
    _run(p, "目  录", 22, bold=True, color=WHITE)
    col1 = C.CONTENTS[:5]
    col2 = C.CONTENTS[5:]
    for ci, col in enumerate((col1, col2)):
        base = ci * 5
        tf = textbox(s, Cm(12.6) + Cm(10.4) * ci, Cm(2.6), Cm(10.0), Cm(14))
        first = True
        for k, name in enumerate(col):
            no = base + k + 1
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            _run(p, f"0{no}  ", 20, bold=True, color=GOLD)
            _run(p, name, 15, bold=True, color=NAVY)
            p.space_after = Pt(20)


def section_divider(prs, no, name):
    s = blank(prs)
    rect(s, Cm(0), Cm(0), SW, SH, NAVY)
    rect(s, Cm(3.5), Cm(8.0), Cm(0.18), Cm(3.0), GOLD)
    tf = textbox(s, Cm(4.2), Cm(6.6), Cm(8), Cm(3.6))
    _run(tf.paragraphs[0], f"0{no}", 60, bold=True, color=GOLD)
    tf = textbox(s, Cm(4.4), Cm(10.4), Cm(24), Cm(2))
    _run(tf.paragraphs[0], name, 26, bold=True, color=WHITE)


def build(prs):
    cover(prs)
    contents(prs)

    # ---- 01 概述与核心定位 ----
    section_divider(prs, 1, C.CONTENTS[0])
    s = blank(prs); content_header(s, 1, C.CONTENTS[0], "项目四大基础优势")
    card_grid(s, C.FOUR_ADVANTAGES, Cm(3.7), cols=2, row_h=Cm(3.35))
    s = blank(prs); content_header(s, 1, C.CONTENTS[0], "四大核心关键梳理")
    card_grid(s, C.FOUR_CORE, Cm(3.7), cols=2, row_h=Cm(3.35))
    s = blank(prs); content_header(s, 1, C.CONTENTS[0], "项目四大运营抓手")
    card_grid(s, C.FOUR_HANDLES, Cm(3.7), cols=2, row_h=Cm(3.35), tinted=True)
    s = blank(prs); content_header(s, 1, C.CONTENTS[0], "项目定位及更名建议")
    rect(s, Cm(1.4), Cm(3.9), Cm(31.0), Cm(3.2), LIGHTGOLD)
    tf = textbox(s, Cm(1.4), Cm(3.9), Cm(31.0), Cm(3.2), anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, C.POSITIONING["更名建议"], 20, bold=True, color=NAVY)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(8)
    _run(p2, "定位口号：" + C.POSITIONING["定位口号"], 16, bold=True, color=GOLD)
    add_bullets(s, [C.POSITIONING["核心定位"], C.POSITIONING["价值主张"]],
                Cm(1.4), Cm(8.0), Cm(31), Cm(6), size=17, gap=18)

    # ---- 02 战略价值 ----
    section_divider(prs, 2, C.CONTENTS[1])
    s = blank(prs); content_header(s, 2, C.CONTENTS[1], "五重赋能叠加驱动未来发展")
    card_grid(s, C.FIVE_EMPOWER, Cm(3.7), cols=2, row_h=Cm(2.55), body_size=11)
    s = blank(prs); content_header(s, 2, C.CONTENTS[1], "对接广州“十五五”战略方向")
    card_grid(s, C.FIFTEEN_FIVE, Cm(3.7), cols=2, row_h=Cm(3.35))
    s = blank(prs); content_header(s, 2, C.CONTENTS[1], "广州出海五大趋势（2026-2030）")
    cards = [(t, "") for t in C.FIVE_TRENDS]
    n = len(C.FIVE_TRENDS)
    cw = Cm(5.9)
    for i, t in enumerate(C.FIVE_TRENDS):
        l = Emu(int(Cm(1.6) + i * (cw + Cm(0.4))))
        rect(s, l, Cm(6.2), cw, Cm(5.2), LIGHT)
        rect(s, l, Cm(6.2), cw, Cm(1.1), NAVY)
        tf = textbox(s, l, Cm(6.2), cw, Cm(1.1), anchor=MSO_ANCHOR.MIDDLE)
        pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        _run(pp, f"趋势 {i+1}", 13, bold=True, color=GOLD)
        tf = textbox(s, l, Cm(7.6), cw, Cm(3.6), anchor=MSO_ANCHOR.MIDDLE)
        pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        _run(pp, t, 15, bold=True, color=NAVY)

    # ---- 03 功能体系与空间布局 ----
    section_divider(prs, 3, C.CONTENTS[2])
    s = blank(prs); content_header(s, 3, C.CONTENTS[2], "五大核心功能")
    card_grid(s, C.FIVE_FUNCTIONS[:4], Cm(3.7), cols=2, row_h=Cm(2.6), body_size=11)
    # 第五个功能单独一行
    box = rect(s, Cm(1.4), Cm(9.55), Cm(31.0), Cm(2.55), LIGHT)
    rect(s, Cm(1.4), Cm(9.55), Cm(0.16), Cm(2.55), GOLD)
    tf = _tf(box); tf.margin_left = Pt(12); tf.margin_top = Pt(9)
    _run(tf.paragraphs[0], C.FIVE_FUNCTIONS[4][0], 15, bold=True, color=NAVY)
    pb = tf.add_paragraph(); _run(pb, C.FIVE_FUNCTIONS[4][1], 11, color=GRAY)
    s = blank(prs); content_header(s, 3, C.CONTENTS[2], "楼层功能定位与业态布局")
    make_table(s, ["区位楼层", "功能定位", "核心业态 / 服务", "数据支撑与参考案例"],
               C.FLOOR_LAYOUT, Cm(1.4), Cm(3.6), Cm(31.0), [3.1, 3.6, 6.4, 5.0],
               header_size=10.5, body_size=8.2, row_h=Cm(1.18), header_h=Cm(0.8))
    s = blank(prs); content_header(s, 3, C.CONTENTS[2], "项目前瞻")
    rect(s, Cm(1.4), Cm(4.4), Cm(31.0), Cm(4.2), LIGHTGOLD)
    tf = textbox(s, Cm(2.2), Cm(4.4), Cm(29.4), Cm(4.2), anchor=MSO_ANCHOR.MIDDLE)
    _run(tf.paragraphs[0], C.PROJECT_SCALE, 17, bold=True, color=NAVY)
    tf = textbox(s, Cm(1.4), Cm(9.6), Cm(31), Cm(2))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, C.SUB_SLOGAN, 16, bold=True, color=GOLD)

    # ---- 04 招商策略 ----
    section_divider(prs, 4, C.CONTENTS[3])
    s = blank(prs); content_header(s, 4, C.CONTENTS[3], "大数据精准锁定客商（一）")
    card_grid(s, C.INVESTMENT[:4], Cm(3.7), cols=2, row_h=Cm(3.35))
    s = blank(prs); content_header(s, 4, C.CONTENTS[3], "大数据精准锁定客商（二）")
    card_grid(s, C.INVESTMENT[4:], Cm(3.7), cols=1, row_h=Cm(2.55), body_size=12)

    # ---- 05 出口体系 ----
    section_divider(prs, 5, C.CONTENTS[4])
    s = blank(prs); content_header(s, 5, C.CONTENTS[4], "出口体系矩阵构建")
    card_grid(s, C.EXPORT_SYSTEM, Cm(3.7), cols=1, row_h=Cm(2.15), body_size=11.5)
    s = blank(prs); content_header(s, 5, C.CONTENTS[4], "六大赋能服务价值")
    card_grid(s, C.SIX_SERVICES, Cm(3.7), cols=2, row_h=Cm(3.35))

    # ---- 06 预期效益 ----
    section_divider(prs, 6, C.CONTENTS[5])
    s = blank(prs); content_header(s, 6, C.CONTENTS[5], "四项核心指标与内外效益")
    card_grid(s, C.BENEFITS, Cm(3.7), cols=2, row_h=Cm(3.35))
    s = blank(prs); content_header(s, 6, C.CONTENTS[5], "未来品牌价值及社会贡献")
    card_grid(s, C.BRAND_VALUE, Cm(3.7), cols=1, row_h=Cm(2.6), body_size=13)

    # ---- 07 运营与合作 ----
    section_divider(prs, 7, C.CONTENTS[6])
    s = blank(prs); content_header(s, 7, C.CONTENTS[6], "收益规划与运营方案")
    card_grid(s, C.OPERATION, Cm(3.7), cols=2, row_h=Cm(2.55), body_size=11)

    # ---- 08 支持保障 ----
    section_divider(prs, 8, C.CONTENTS[7])
    s = blank(prs); content_header(s, 8, C.CONTENTS[7], "黄埔区政府渠道与政策支持")
    card_grid(s, C.GOV_SUPPORT, Cm(3.7), cols=2, row_h=Cm(2.55), body_size=10.5)
    s = blank(prs); content_header(s, 8, C.CONTENTS[7], "中广协文创与IP专委会资源支持")
    card_grid(s, C.CGC_RESOURCES, Cm(3.7), cols=2, row_h=Cm(3.35), accent=GOLD)

    # ---- 09 总结与未来方向 ----
    section_divider(prs, 9, C.CONTENTS[8])
    s = blank(prs); content_header(s, 9, C.CONTENTS[8], "项目总结")
    add_bullets(s, C.SUMMARY_POINTS, Cm(1.4), Cm(3.7), Cm(31), Cm(14),
                size=14, gap=13)
    # 未来方向(重点)
    s = blank(prs); content_header(s, 9, C.CONTENTS[8], C.FUTURE_TITLE)
    rect(s, Cm(1.4), Cm(3.6), Cm(31.0), Cm(2.2), LIGHTGOLD)
    tf = textbox(s, Cm(1.7), Cm(3.6), Cm(30.4), Cm(2.2), anchor=MSO_ANCHOR.MIDDLE)
    _run(tf.paragraphs[0], C.FUTURE_INTRO, 12.5, color=NAVY)
    card_grid(s, C.FUTURE_POINTS, Cm(6.2), cols=2, row_h=Cm(3.0),
              body_size=11, accent=GOLD)

    # ---- 附件表 ----
    s = blank(prs); content_header(s, 9, "附件", "附件一：2025年广州出口TOP20品类")
    make_table(s, ["排名", "品类", "代表品牌", "广州口岸交易额", "核心说明"],
               C.EXPORT_TOP20, Cm(1.4), Cm(3.5), Cm(31.0), [1.3, 3.6, 3.6, 3.0, 6.5],
               header_size=10, body_size=7.6, row_h=Cm(0.66), header_h=Cm(0.72))
    tf = textbox(s, Cm(1.4), Cm(18.05), Cm(31), Cm(0.9))
    _run(tf.paragraphs[0], C.EXPORT_TOP20_NOTE, 8.5, color=GRAY)

    s = blank(prs); content_header(s, 9, "附件", "附件二：2025年广州消费类出口20强")
    make_table(s, ["排名", "品类", "出口额(亿元)", "同比", "核心出口品牌", "核心市场"],
               C.CONSUMER_TOP20, Cm(1.4), Cm(3.5), Cm(31.0),
               [1.3, 3.5, 2.6, 2.0, 4.2, 3.2],
               header_size=9.5, body_size=7.6, row_h=Cm(0.66), header_h=Cm(0.72))
    tf = textbox(s, Cm(1.4), Cm(18.05), Cm(31), Cm(0.9))
    _run(tf.paragraphs[0], C.CONSUMER_TOP20_NOTE, 8.5, color=GRAY)

    # ---- 封底 ----
    s = blank(prs)
    rect(s, Cm(0), Cm(0), SW, SH, NAVY)
    rect(s, Cm(11.9), Cm(10.7), Cm(10), Cm(0.1), GOLD)
    tf = textbox(s, Cm(2), Cm(6.6), Cm(29.8), Cm(3))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, "谢谢观看", 40, bold=True, color=WHITE)
    tf = textbox(s, Cm(2), Cm(11.4), Cm(29.8), Cm(2.5))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, f"{C.OWNER_ORG}   {C.OWNER_PERSON}", 15, bold=True, color=GOLD)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(6)
    _run(p2, C.DOC_DATE, 12, color=LIGHT)


def main():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    build(prs)
    os.makedirs(OUT_DIR, exist_ok=True)
    out = os.path.join(OUT_DIR, "广州黄埔九佛TOD全球自贸365街区_策划方案.pptx")
    prs.save(out)
    print("saved:", out, "slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
