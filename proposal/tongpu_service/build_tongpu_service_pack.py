# -*- coding: utf-8 -*-
"""同浦汇 · 30场活动与科技企业服务中心筹备 — 复核定稿包

交付四件套：
1. 执行台账.xlsx
2. 业务承接策划案.docx
3. 业务承接策划案.pptx
4. 业务承接合作协议.docx（可直接签署）
"""
from __future__ import annotations

import shutil
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.shared import Cm, Pt, RGBColor as DocRGB
from openpyxl import Workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from openpyxl.worksheet.datavalidation import DataValidation
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt as PPt

HERE = Path(__file__).resolve().parent
OUT = HERE.parent
ART = Path("/opt/cursor/artifacts")

# —— 统一口径 ——
PARTY_A = "上海同浦汇科技有限公司"
PARTY_A_SHORT = "同浦汇"
PARTY_B = "上海市杨浦区科技企业服务中心"
PARTY_B_SHORT = "服中心"
PARTY_C = "上海市杨浦区科技企业联合会"
OWNER = "杨浦科创集团"
PARK = "创智汇"
ACADEMIC = "复旦大学住房政策研究中心"
CONTACT_TPH = "高辰辰　15339617481　产业招服"
PERIOD = "2026.08—2027.07"
DATE_CN = "2026年9月1日"

EVENTS = [
    ("E1", "2026.08", "启动", "赛道切换说明会 · 园区开放日", "把新方向讲清楚，完成主业匹配叙事"),
    ("E2", "2026.08", "启动", "中介渠道推广日（建造 / 园区）", "点火带客，不谈租金对赌"),
    ("E3", "2026.08", "启动", "五大行与资管对接日", "面对施工企业支付能力弱的现实，先找金融出口"),
    ("E4", "2026.08", "启动", "政府部门与载体沟通日", "确认「科技赋能建筑」政策包装"),
    ("E5", "2026.08", "启动", "进博绿色建造成果承接交流日", "借势施工机器人 + 绿色低碳，不替代年包出海"),
    ("A1", "2026.09", "智能建造", "施工机器人与智能装备一日营", "3F能力线：装备可见、可体验"),
    ("C1", "2026.09", "出海准备", "东南亚模块化建筑案例沙龙（国内场）", "讲案例与供应链，不含领事"),
    ("B1", "2026.09", "政策", "智能建造企业高新认定路径沙龙", "经营范围写入智能施工 / 绿色节能"),
    ("A2", "2026.10", "智能建造", "模块化 / 装配式建筑工作坊", "可出口产品与就地建厂逻辑"),
    ("D1", "2026.10", "展示", "500㎡出海展厅共创工作坊", "把展陈从概念落到平面"),
    ("B2", "2026.10", "政策", "创新券与数字化设计工具沙龙", "服中心作为服务机构的入口"),
    ("A3", "2026.11", "智能建造", "绿色低碳与零碳工地实务营", "对接进博观察与越秀研学思路"),
    ("B3", "2026.11", "政策", "YOUNG立方 · 智能建造政策沙龙", "活动经费申报的合规场次"),
    ("A4", "2026.12", "智能建造", "智能测绘与数字设计对接", "同济设计资源协同，不替代复旦住房政策线"),
    ("D2", "2026.12", "展示", "行业白皮书开题沙龙", "联合复旦形成可传播文本"),
    ("F1", "2026.12", "收官", "智能建造年度 Demo Day", "年度成果可核验"),
    ("A5", "2027.01", "智能建造", "建筑机器人场景体验日", "3F体验与招商转化"),
    ("C2", "2027.01", "出海准备", "高校成果转化 · 智能建造", "复旦住房政策 + 同济设计双校协同"),
    ("B4", "2027.02", "政策", "高企认定冲刺（智能施工专场）", "企业自行申报，服中心辅导不承诺获批"),
    ("D3", "2027.02", "展示", "零碳体验馆研学方案工作坊", "中小学生研学运营模型"),
    ("A6", "2027.03", "智能建造", "建造科技黑客松（OPC超级个体）", "保留科技社群带流，服务于新赛道"),
    ("C3", "2027.03", "出海准备", "华东智能设计产业集群对接", "华东产能与设计优势"),
    ("F2", "2027.03", "收官", "通往海外订单季度圆桌（国内）", "订单逻辑研讨，领事另计价"),
    ("C4", "2027.04", "出海准备", "华南制造基地联动对接", "华东设计 × 华南制造"),
    ("D4", "2027.04", "展示", "新材料与绿色建材联展", "5F展陈轮换"),
    ("C5", "2027.05", "出海准备", "专精特新 · 建造科技路演", "申报奖项节奏的企业侧入口"),
    ("D5", "2027.05", "展示", "展厅内容迭代与市集体验日", "人气场，≤30人核验口径"),
    ("C6", "2027.06", "出海准备", "一带一路目标国产业地图沙龙", "城市化率<50%国家研究场"),
    ("A7", "2027.06", "智能建造", "智能建造与智能家居交叉应用营", "产品出口与场景验证"),
    ("F3", "2027.07", "收官", "年度收官 · 下一年度预热圆桌", "验收30场 + 预热下一年度"),
]
assert len(EVENTS) == 30

# —— Excel ——
THIN = Border(
    left=Side(style="thin", color="D6D0E8"),
    right=Side(style="thin", color="D6D0E8"),
    top=Side(style="thin", color="D6D0E8"),
    bottom=Side(style="thin", color="D6D0E8"),
)
FILL_H = PatternFill("solid", fgColor="3A206E")
FILL_T = PatternFill("solid", fgColor="1A1140")
FILL_A = PatternFill("solid", fgColor="F7F4FF")
FILL_B = PatternFill("solid", fgColor="FFFFFF")
FILL_NOTE = PatternFill("solid", fgColor="FFF8E7")
FONT_H = Font(name="微软雅黑", bold=True, color="FFFFFF", size=11)
FONT_T = Font(name="微软雅黑", bold=True, color="FFFFFF", size=14)
FONT_N = Font(name="微软雅黑", size=10, color="2A2240")
FONT_M = Font(name="微软雅黑", size=10, color="5A4E7A")
AL_C = Alignment(horizontal="center", vertical="center", wrap_text=True)
AL_L = Alignment(horizontal="left", vertical="center", wrap_text=True)


def _ws_title(ws, title, subtitle="", cols=8):
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=cols)
    c = ws.cell(1, 1, title)
    c.font = FONT_T
    c.fill = FILL_T
    c.alignment = AL_L
    ws.row_dimensions[1].height = 28
    if subtitle:
        ws.merge_cells(start_row=2, start_column=1, end_row=2, end_column=cols)
        c2 = ws.cell(2, 1, subtitle)
        c2.font = FONT_M
        c2.fill = FILL_NOTE
        c2.alignment = AL_L
        ws.row_dimensions[2].height = 22


def _hdr(ws, row, headers):
    for i, h in enumerate(headers, 1):
        cell = ws.cell(row, i, h)
        cell.font = FONT_H
        cell.fill = FILL_H
        cell.alignment = AL_C
        cell.border = THIN


def _cell(ws, r, c, v, center=False):
    cell = ws.cell(r, c, v)
    cell.font = FONT_N
    cell.alignment = AL_C if center else AL_L
    cell.border = THIN
    cell.fill = FILL_A if r % 2 == 0 else FILL_B
    return cell


def _widths(ws, widths):
    for i, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = w


def build_xlsx() -> Path:
    wb = Workbook()
    ws = wb.active
    ws.title = "00-封面说明"
    _ws_title(ws, f"{PARK}30场活动与科技企业服务中心筹备　·　业务承接合作策划案",
              f"提交 {PARTY_A_SHORT}　｜　提出方 {PARTY_B_SHORT}　｜　{DATE_CN}　｜　周期 {PERIOD}")
    rows = [
        ("字段", "内容"),
        ("提交对象", PARTY_A),
        ("提出方 / 执行方", PARTY_B),
        ("联合协办", PARTY_C),
        ("学术支持", ACADEMIC),
        ("同浦汇联系人", CONTACT_TPH),
        ("业主", f"{OWNER}（杨浦国资）"),
        ("园区 / 项目", f"{PARK}（五角场·创智天地片区）"),
        ("新赛道", "智能建造 × 建筑产业出海"),
        ("政策包装", "现代服务业 + 新一代信息技术（科技赋能建筑）"),
        ("工作包 A", "30场活动全案执行（30万/年；每场≤30人；负责人≤30%）"),
        ("工作包 B", "科技企业服务中心90天筹备（筹备不另向同浦汇收费）"),
        ("活动内部分成", "服中心70% / 同浦汇30%（建议，待签约确认）"),
        ("政策分成", "同浦汇38% / 服中心62%（锁版）"),
        ("付款", "签约后7日内50%启动款；Q3、Q4、次年Q1、次年Q2各10%；年终收官10%尾款"),
        ("机密", "机密 · 仅供同浦汇合作沟通，不作为对园区/投资方的承诺函"),
    ]
    for i, (a, b) in enumerate(rows, 4):
        _cell(ws, i, 1, a, center=True)
        _cell(ws, i, 2, b)
        if i == 4:
            ws.cell(i, 1).font = FONT_H
            ws.cell(i, 1).fill = FILL_H
            ws.cell(i, 2).font = FONT_H
            ws.cell(i, 2).fill = FILL_H
    ws.cell(21, 1, "工作表索引").font = Font(name="微软雅黑", bold=True, size=12, color="3A206E")
    _hdr(ws, 22, ["工作表", "用途", "对应文件"])
    for i, row in enumerate([
        ("00-封面说明", "口径、主体、分成总览", "PPT + Word + 协议"),
        ("01-昨日共识", "8月31日交流要点，作为共同事实", "PPT + Word"),
        ("02-分工矩阵", "同浦汇 / 服中心 / 园区 / 复旦 / 业主", "PPT + Word + 协议"),
        ("03-30场总表", "30场排期、线条、状态（可勾选）", "PPT + Word + 协议附件"),
        ("04-单场标准", "人数、负责人、节奏、转化、加购", "PPT + Word + 协议"),
        ("05-服中心90天", "挂牌筹备甘特与前置条件", "PPT + Word + 协议"),
        ("06-商务付款", "30万拆分到季度与内部结算", "PPT + Word + 协议"),
        ("07-政策分成", "上限测算与38/62", "PPT + Word + 协议"),
        ("08-边界风险", "不承诺清单 + 审核三句口径", "PPT + Word + 协议"),
        ("09-下一步", "请同浦汇确认的五件事", "PPT + Word"),
    ], 23):
        for j, v in enumerate(row, 1):
            _cell(ws, i, j, v)
    _widths(ws, [22, 55, 28])

    ws = wb.create_sheet("01-昨日共识")
    _ws_title(ws, "昨日共识（2026-08-31）", "创智汇赛道调整交流 · 作为双方共同事实，不另做预测承诺", 3)
    _hdr(ws, 3, ["序号", "类型", "要点"])
    consensus = [
        ("1", "赛道", "创智汇上周已确定更换赛道：由原AI+IP内容线，转向与建筑设计、智能建造结合的方向，并与同济设计资源、中建相关方主业相匹配。"),
        ("2", "赛道", "新方向按「现代服务业 + 新一代信息技术」包装，科技赋能建筑，不与高新/载体政策申报发生本质冲突。"),
        ("3", "目标", "核心目标不是把地产施工企业招来租办公室，而是搭建「建筑产业出海」集群：设计、测绘、模块化、绿色低碳、施工机器人可成为出口内容。"),
        ("4", "风险", "会上研判：传统建筑施工景气与回款能力明显弱于人工智能赛道，若仍按招租逻辑推进，后续难以说服审核方，也难以支撑活动转化。"),
        ("5", "口径", "对外口径须谨慎：业主写杨浦科创集团；不宜简单用「中建」名义对接政府资源。"),
        ("6", "资源", f"{ACADEMIC}可提供行业数据、活动组织、专家网络与海外专题经验，作为白皮书、授牌与出海研究的学术支持，不替代同浦汇的园区接口地位。"),
        ("7", "出海", "出海优先城市化率较低的一带一路、东南亚、中亚、非洲、中东市场；产品抓手为绿色低碳 + 建筑施工机器人 + 模块化建筑。"),
        ("8", "政策", "引导入驻企业把经营范围写入智能施工、绿色节能，自行完成高新技术企业认定；可参考零碳体验馆研学模式，不完全依赖补贴。"),
        ("9", "空间", "展厅约500㎡做建筑出海科技与模块化产品展示，其余空间承载新材料等。"),
        ("10", "领馆", "驻沪领馆资源按目标国匹配（会上口径约80+），领事到访不纳入30场年包。"),
        ("11", "节奏", "合同口径已调整、运营成本已经发生；需尽快形成可交给审核方工程部门的初步方案。"),
    ]
    for i, row in enumerate(consensus, 4):
        for j, v in enumerate(row, 1):
            _cell(ws, i, j, v, center=(j < 3))
        ws.row_dimensions[i].height = 36
    ws.merge_cells("A16:C16")
    ws.cell(16, 1, "定位切换：AI+IP / 数字内容无界共创港  →  智能建造 × 建筑产业出海　｜　政策包装：现代服务业 + 新一代信息技术")
    ws.cell(16, 1).fill = FILL_NOTE
    ws.cell(16, 1).font = Font(name="微软雅黑", bold=True, size=10, color="3A206E")
    _widths(ws, [8, 10, 90])

    ws = wb.create_sheet("02-分工矩阵")
    _ws_title(ws, "分工矩阵", "同浦汇管关系，服中心管交付；园区管销售促成", 4)
    _hdr(ws, 3, ["主体", "责任", "主责/协同", "是否承接范围内"])
    for i, row in enumerate([
        (PARTY_A_SHORT, "园区接口、在谈客户、带客名单协同、建群与回访台账、向园区转交促成线索", "主责客户关系", "接口保留，不替换"),
        (PARTY_B_SHORT, "30场策划与现场执行、服中心挂牌筹备、政策申报材料、月报验收", "主责执行", "是，工作包A+B"),
        ("科技企业联合会", "联合协办、行业组织协同、活动背书", "协同", "联合协办"),
        ("复旦住房政策研究中心", "学术支持、白皮书、专家与行业数据、海外专题研究（按场次确认）", "按场次确认", "学术支持，不替代接口"),
        (f"{PARK} / 园区", "场地物业、导视、样板间、在谈名单、销售促成", "主责销售", "不纳入承接承诺"),
        (OWNER, "空间业主；重大口径与场地条件确认", "业主确认", "口径与场地条件"),
    ], 4):
        for j, v in enumerate(row, 1):
            _cell(ws, i, j, v)
        ws.row_dimensions[i].height = 32
    _widths(ws, [22, 55, 16, 22])

    ws = wb.create_sheet("03-30场总表")
    _ws_title(ws, "30场活动执行总表（计入年包）",
              "每场≤30人　｜　负责人≤30%　｜　状态可下拉　｜　另计价项目见「04-单场标准」", 14)
    _hdr(ws, 3, ["序号", "编号", "月份", "线条", "活动名称", "本场作用", "建议周次",
                 "人数上限", "负责人占比上限", "主责", "同浦汇协同", "状态", "线索条数", "备注"])
    for i, (code, month, line, name, role) in enumerate(EVENTS, 1):
        r = i + 3
        vals = [i, code, month, line, name, role, "", 30, "≤30%", "服中心执行", "带客/回访", "未启动", "", ""]
        for j, v in enumerate(vals, 1):
            _cell(ws, r, j, v, center=j in (1, 2, 3, 7, 8, 9, 12, 13))
        ws.row_dimensions[r].height = 28
    dv = DataValidation(type="list", formula1='"未启动,策划中,邀约中,已执行,已复盘,延期,取消"', allow_blank=True)
    ws.add_data_validation(dv)
    dv.add("L4:L33")
    ws.merge_cells("A35:N35")
    ws.cell(35, 1, "场次合计：30　｜　启动5　智能建造7　政策4　出海准备6　展示5　收官3　｜　另计价不计入本表")
    ws.cell(35, 1).fill = FILL_NOTE
    ws.cell(35, 1).font = Font(name="微软雅黑", bold=True, size=10, color="3A206E")
    _widths(ws, [6, 8, 10, 10, 36, 32, 10, 10, 14, 12, 12, 10, 10, 16])

    ws = wb.create_sheet("04-单场标准")
    _ws_title(ws, "单场标准 / 转化 / 加购",
              "T-14策划邀约 → T-7确认名单 → T-1场地物料 → D日执行 → T+1~7建群/回访摘要", 3)
    _hdr(ws, 3, ["项目", "口径", "说明"])
    for i, row in enumerate([
        ("人数", "每场≤30人（可核验）", "签到表核验"),
        ("负责人占比", "≤30%，不做100%承兑", "覆盖早期约50%口径；本版统一≤30%"),
        ("全年触达", "全年约600+人次", "30×约20人量级，不写千人场"),
        ("执行节奏", "T-14→T-7→T-1→D日→T+1~7", "档期灵活，不是一个月一场"),
        ("转化闭环", "到场→建群→看场→同浦汇回访→园区销售促成→复盘回流", "活动不对租金去化对赌"),
        ("交付物", "方案、签到、现场执行、月报；回访只交摘要知会", "不要求园区逐户回访"),
    ], 4):
        for j, v in enumerate(row, 1):
            _cell(ws, i, j, v)
    ws.cell(11, 1, "另计价（不进30万包）").font = Font(name="微软雅黑", bold=True, size=11, color="3A206E")
    _hdr(ws, 12, ["项目", "内容", "计价原则"])
    for i, row in enumerate([
        ("出海专题执行", "境外对接、订单撮合、海外考察", "10–20万/场次或专项，一事一议"),
        ("领事到访接待", "驻沪领馆匹配目标国", "按场次另计价，与挂牌分开"),
        ("挂牌仪式", "服中心 / 复旦战略合作授牌", "3万/个起，挂牌≠领事活动"),
        ("战略签约大会升级", "媒体、会场、嘉宾规格升级", "在预热沙龙之外另报价"),
    ], 13):
        for j, v in enumerate(row, 1):
            _cell(ws, i, j, v)
    _widths(ws, [18, 48, 36])

    ws = wb.create_sheet("05-服中心90天")
    _ws_title(ws, "科技企业服务中心筹备 · 90天", "挂牌仪式另计价　｜　筹备不另向同浦汇收费", 8)
    _hdr(ws, 3, ["阶段", "窗口", "主题", "交付", "W1-2", "W3-6", "W7-10", "W11-12"])
    for i, row in enumerate([
        ("1", "D0–15", "叙事与主体", "统一新赛道一页纸；确认挂牌主体、2人配置、创智汇接待位/500㎡展区界面", "●", "", "", ""),
        ("2", "D16–45", "制度与申报", "服务制度、辅导SOP、载体认定材料包、创新券机构申请、高新路径清单", "", "●", "", ""),
        ("3", "D46–75", "产品化", "展陈大纲、零碳研学模型、白皮书开题、首批可辅导企业名单", "", "", "●", ""),
        ("4", "D76–90", "可过审包", "给审核方工程部门的初步方案、挂牌仪式报价（另计价）、Q4场次点火", "", "", "", "●"),
    ], 4):
        for j, v in enumerate(row, 1):
            _cell(ws, i, j, v, center=j != 4)
        ws.row_dimensions[i].height = 40
    ws.cell(9, 1, "前置条件（政策）").font = Font(name="微软雅黑", bold=True, size=11, color="3A206E")
    ws.merge_cells("A10:H10")
    ws.cell(10, 1, "①项目成为经认定载体 / 成果转化服务平台　②服中心作为创新券服务机构")
    ws.merge_cells("A12:H12")
    ws.cell(12, 1, "空间：创智汇（一期）约6600㎡　｜　3F约2850㎡孵化器+办公　｜　5F约3670㎡展厅+贸易　｜　专题约500㎡建筑产业出海主题展示　｜　办公约3.3元/㎡/天　｜　物业约13.8元/㎡/月")
    ws.cell(12, 1).fill = FILL_NOTE
    _widths(ws, [8, 10, 12, 55, 8, 8, 8, 10])

    ws = wb.create_sheet("06-商务付款")
    _ws_title(ws, "商务付款与内部结算", "对园区一口价　｜　内部两套分成互不混用", 6)
    ws.cell(3, 1, "对园区 · 活动年包30万元拆分").font = Font(name="微软雅黑", bold=True, size=11, color="3A206E")
    _hdr(ws, 4, ["节点", "比例", "金额（万元）", "累计（万元）", "挂钩条件", "备注"])
    for i, row in enumerate([
        ("签约后7日内 · 启动款", "50%", 15, 15, "启动策划与前8场", "到账后T+14交执行手册"),
        ("2026 Q3", "10%", 3, 18, "Q3场次完成度 + 月报", "8–9月为主"),
        ("2026 Q4", "10%", 3, 21, "Q4场次完成度 + 月报", "10–12月"),
        ("2027 Q1", "10%", 3, 24, "Q1场次完成度 + 月报", "1–3月"),
        ("2027 Q2", "10%", 3, 27, "Q2场次完成度 + 月报", "4–6月；7月收官计入年终"),
        ("年终收官尾款", "10%", 3, 30, "30场核验 + 年终月报", "完成后15日内"),
        ("合计", "100%", 30, 30, "", ""),
    ], 5):
        for j, v in enumerate(row, 1):
            _cell(ws, i, j, v, center=j > 1)
    ws.cell(13, 1, "内部结算（建议，签约确认）").font = Font(name="微软雅黑", bold=True, size=11, color="3A206E")
    _hdr(ws, 14, ["收入类型", "基数（万元）", "服中心", "同浦汇", "服中心金额", "同浦汇金额"])
    for i, row in enumerate([
        ("活动年包", 30, "70%", "30%", 21, 9),
        ("政策收益（到账后，示例不作承诺）", "待到账", "62%", "38%", "待到账×62%", "待到账×38%"),
        ("服中心筹备", 0, "不另收费", "—", 0, 0),
    ], 15):
        for j, v in enumerate(row, 1):
            _cell(ws, i, j, v, center=j > 1)
    ws.merge_cells("A19:F19")
    ws.cell(19, 1, "招商佣金：2个月净租金（首年不重复），由园区销售闭环触发，不计入上表。门票/赞助可冲抵活动成本，不替代30万打包价。")
    ws.cell(19, 1).fill = FILL_NOTE
    _widths(ws, [32, 14, 12, 12, 16, 28])

    ws = wb.create_sheet("07-政策分成")
    _ws_title(ws, "政策工具箱与分成（上限测算，非保证获批）",
              "①项目成为经认定载体/成果转化服务平台　②服中心作为创新券服务机构", 6)
    _hdr(ws, 3, ["项目", "年上限（万元）", "十年（万元）", "同浦汇38%", "服中心62%", "备注"])
    for i, row in enumerate([
        ("载体（平台100+基地10）", 110, 1100, 418, 682, "认定后才可申报"),
        ("活动 YOUNG立方封顶", 200, 2000, 760, 1240, "投入×50%，年包30万大约对应约15万量级"),
        ("申报奖项（培育节奏测算）", "80→530", 3820, 1451.6, 2368.4, "高企/专精特新/小巨人等，按到账"),
        ("三部分合计上限", "390→840", 6920, 2629.6, 4290.4, "上限不是保底"),
    ], 4):
        for j, v in enumerate(row, 1):
            _cell(ws, i, j, v, center=j > 1)
    ws.merge_cells("A9:F9")
    ws.cell(9, 1, "服务业引导资金≤300万为项目制，未计入十年年表。成功费以资金到账为准，不向企业承诺必然获批。")
    ws.cell(9, 1).fill = FILL_NOTE
    _widths(ws, [28, 16, 14, 14, 14, 40])

    ws = wb.create_sheet("08-边界风险")
    _ws_title(ws, "不承诺清单与审核口径", "转交审核方时建议原样保留三句话", 2)
    _hdr(ws, 3, ["类型", "条文"])
    for i, row in enumerate([
        ("不承接 / 不承诺", "不取代同浦汇作为园区产业招服入口与客户关系主人"),
        ("不承接 / 不承诺", "不做租金去化对赌，也不做「必须租出去」的必要性条款（办公约3.3元/㎡/天高于周边）"),
        ("不承接 / 不承诺", "不承诺必带外资企业、不承诺领事必到场"),
        ("不承接 / 不承诺", "不把出海专题、领事到访、挂牌仪式装进30万元年包"),
        ("不承接 / 不承诺", "不向企业承诺高新/载体/补贴必然获批"),
        ("不承接 / 不承诺", "园区销售促成仍由园区负责；我方与同浦汇负责带到场、建群、回访摘要"),
        ("审核口径1", "我们不是在招「付不起租金的施工队来租办公室」。"),
        ("审核口径2", "我们是在用科技企业服务中心，把智能建造产品、模块化建筑和绿色低碳能力组织成可出海的集群。"),
        ("审核口径3", "业主是杨浦科创集团；政府对接使用科技与产业服务口径，不以「中建」名义包装。"),
        ("KPI场次", "12个月内完成30场可核验活动"),
        ("KPI人数", "单场签到≤30人，全年触达约600+"),
        ("KPI月报", "每月1份，含名单摘要、照片、线索条数"),
        ("KPI服中心", "90天内完成挂牌材料与制度包，具备申报条件"),
        ("KPI转化", "每场输出看场线索，由同浦汇回访、园区促成；不设租金对赌"),
    ], 4):
        for j, v in enumerate(row, 1):
            _cell(ws, i, j, v)
        ws.row_dimensions[i].height = 28
    _widths(ws, [18, 90])

    ws = wb.create_sheet("09-下一步")
    _ws_title(ws, "请同浦汇确认", "五件事确认后，启动款到账14日内交付细化执行手册", 4)
    _hdr(ws, 3, ["序号", "确认事项", "状态", "同浦汇意见"])
    for i, row in enumerate([
        ("1", "同浦汇确认：承接范围（30场 + 服中心筹备）与不承接边界无异议", "待确认", ""),
        ("2", "确认内部结算：活动年包70/30、政策收益38/62", "待确认", ""),
        ("3", "确认8–9月前8场档期与场地", "待确认", ""),
        ("4", "确认向审核方提交的业主口径（杨浦科创集团）与赛道叙事", "待确认", ""),
        ("5", "签约后7日内启动款到位，T+14交付细化执行手册", "待确认", ""),
    ], 4):
        for j, v in enumerate(row, 1):
            _cell(ws, i, j, v, center=j in (1, 3))
        ws.row_dimensions[i].height = 28
    ws.merge_cells("A10:D10")
    ws.cell(10, 1, f"联系人：{CONTACT_TPH}　｜　提出方：{PARTY_B}　｜　{DATE_CN}")
    ws.cell(10, 1).fill = FILL_NOTE
    _widths(ws, [8, 70, 12, 20])

    path = OUT / f"{PARTY_A_SHORT}_30场活动与科技企业服务中心筹备_执行台账.xlsx"
    wb.save(path)
    return path


# —— Word helpers ——
SONG, HEI = "宋体", "黑体"


def _set_run_font(run, font_name, size, bold=False):
    run.font.name = font_name
    run.font.size = Pt(size)
    run.bold = bold
    rpr = run._element.get_or_add_rPr()
    rfonts = rpr.find(qn("w:rFonts"))
    if rfonts is None:
        rfonts = rpr.makeelement(qn("w:rFonts"), {})
        rpr.append(rfonts)
    rfonts.set(qn("w:eastAsia"), font_name)
    rfonts.set(qn("w:ascii"), font_name)
    rfonts.set(qn("w:hAnsi"), font_name)


def wp(doc, text="", *, font=SONG, size=12, bold=False, align=WD_ALIGN_PARAGRAPH.LEFT,
       first=None, after=6, before=0):
    para = doc.add_paragraph()
    para.alignment = align
    pf = para.paragraph_format
    pf.space_after = Pt(after)
    pf.space_before = Pt(before)
    pf.line_spacing = 1.5
    if first is not None:
        pf.first_line_indent = Cm(first)
    if text:
        run = para.add_run(text)
        _set_run_font(run, font, size, bold=bold)
    return para


def wh(doc, text):
    return wp(doc, text, font=HEI, size=14, bold=True, before=12, after=8)


def wbod(doc, text):
    return wp(doc, text, size=12, align=WD_ALIGN_PARAGRAPH.JUSTIFY, first=0.74, after=6)


def wbullet(doc, text):
    para = doc.add_paragraph(style="List Bullet")
    para.paragraph_format.line_spacing = 1.4
    para.paragraph_format.space_after = Pt(3)
    run = para.add_run(text)
    _set_run_font(run, SONG, 12)
    return para


def wtable(doc, rows, col_widths=None):
    table = doc.add_table(rows=len(rows), cols=len(rows[0]))
    table.style = "Table Grid"
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = ""
            p = cell.paragraphs[0]
            run = p.add_run(str(val))
            _set_run_font(run, SONG, 10.5, bold=(ri == 0))
    if col_widths:
        for row in table.rows:
            for i, w in enumerate(col_widths):
                row.cells[i].width = Cm(w)
    doc.add_paragraph()
    return table


def build_docx_plan() -> Path:
    doc = Document()
    for s in doc.sections:
        s.top_margin = Cm(2.4)
        s.bottom_margin = Cm(2.4)
        s.left_margin = Cm(2.6)
        s.right_margin = Cm(2.6)

    wp(doc, "提交对象", font=HEI, size=11, bold=True, align=WD_ALIGN_PARAGRAPH.CENTER, after=2)
    wp(doc, PARTY_A, font=HEI, size=16, bold=True, align=WD_ALIGN_PARAGRAPH.CENTER, after=10)
    wp(doc, f"{PARK}30场活动与科技企业服务中心筹备", font=HEI, size=18, bold=True,
       align=WD_ALIGN_PARAGRAPH.CENTER, after=4)
    wp(doc, "业务承接合作策划案（复核定稿）", font=HEI, size=16, bold=True,
       align=WD_ALIGN_PARAGRAPH.CENTER, after=10)
    wp(doc, "承接范围：30场可核验活动  ＋  科技企业服务中心挂牌筹备",
       align=WD_ALIGN_PARAGRAPH.CENTER, after=2)
    wp(doc, "新赛道：智能建造 × 建筑产业出海", align=WD_ALIGN_PARAGRAPH.CENTER, after=2)
    wp(doc, "政策包装：现代服务业 + 新一代信息技术（科技赋能建筑）",
       align=WD_ALIGN_PARAGRAPH.CENTER, after=8)
    wp(doc, f"提出方：{PARTY_B}", align=WD_ALIGN_PARAGRAPH.CENTER, after=2)
    wp(doc, f"联合：{PARTY_C}", align=WD_ALIGN_PARAGRAPH.CENTER, after=2)
    wp(doc, f"学术支持：{ACADEMIC}", align=WD_ALIGN_PARAGRAPH.CENTER, after=2)
    wp(doc, f"{DATE_CN}　｜　执行周期 {PERIOD}", align=WD_ALIGN_PARAGRAPH.CENTER, after=4)
    wp(doc, "机密 · 仅供同浦汇合作沟通，不作为对园区/投资方的承诺函",
       font=SONG, size=10, align=WD_ALIGN_PARAGRAPH.CENTER, after=16)

    wh(doc, "致同浦汇")
    wbod(doc,
         "2026年8月31日关于创智汇赛道调整的交流已把方向说清楚：园区不再以原AI+IP内容线作为主叙事，"
         "而是转向智能建造与建筑产业出海，并用现代服务业、新一代信息技术的口径完成政策包装。"
         "合同口径已经调整，运营成本已经发生，审核方工程部门需要尽快看到一份可执行、可核验、可过审的答卷。")
    wbod(doc,
         f"据此，{PARTY_B}提出：承接同浦汇在本项目上的两项具体业务——一是把30场活动按新赛道办完；"
         "二是把科技企业服务中心筹备到可挂牌、可申报。同浦汇继续做创智汇面前的产业招服入口和客户关系主人，"
         "不把园区接口让出去。本文件为复核定稿版，便于内部对齐后转交审核材料，并作为签署《业务承接合作协议》的附件说明；"
         "不作为对园区或投资方的单方承诺函。")

    wh(doc, "一、昨日共识：为什么现在要承接")
    wbod(doc, "以下内容按8月31日交流整理，作为双方共同事实，不另做行业预测承诺。")
    for t in [
        "1. 创智汇上周已确定更换赛道：由原AI+IP内容线，转向与建筑设计、智能建造结合的方向，并与同济设计资源、中建相关方主业相匹配。",
        "2. 新方向按「现代服务业 + 新一代信息技术」包装，科技赋能建筑，不与高新/载体政策申报发生本质冲突。",
        "3. 核心目标不是把地产施工企业招来租办公室，而是搭建「建筑产业出海」集群。",
        "4. 会上研判：传统建筑施工景气与回款能力明显弱于人工智能赛道，若仍按招租逻辑推进，后续难以说服审核方。",
        "5. 对外口径须谨慎：业主写杨浦科创集团；不宜简单用「中建」名义对接政府资源。",
        f"6. {ACADEMIC}可提供行业数据、活动组织、专家网络与海外专题经验，不替代同浦汇的园区接口地位。",
        "7. 出海优先城市化率较低的一带一路、东南亚、中亚、非洲、中东市场；产品抓手为绿色低碳 + 施工机器人 + 模块化建筑。",
        "8. 政策路径：引导入驻企业把经营范围写入智能施工、绿色节能，自行完成高新技术企业认定。",
        "9. 空间抓手：展厅约500㎡做建筑出海科技与模块化产品展示。",
        "10. 驻沪领馆资源按目标国匹配；领事到访不纳入30场年包。",
        "11. 合同口径已调整、运营成本已经发生；需尽快形成可交给审核方工程部门的初步方案。",
    ]:
        wbullet(doc, t)
    wbod(doc, "一句话：旧定位「AI+IP / 数字内容无界共创港」让位给新定位「智能建造 × 建筑产业出海」。30场年包的商务结构沿用，主题线与服中心筹备按新赛道重写。")

    wh(doc, "二、承接范围与不承接边界")
    wbod(doc, "2.1 承接什么")
    wtable(doc, [
        ("工作包", "具体内容"),
        ("30场活动全案执行", "策划、档期、邀约协同、现场执行、月报、可核验交付"),
        ("科技企业服务中心筹备", "挂牌路径、制度与人员、载体认定材料、创新券机构申请、高新辅导SOP、展陈与研学大纲"),
        ("新赛道叙事与审核材料", "把昨日共识写成可过审的一页纸 + 90天工作包"),
        ("学术与白皮书协同", f"对接{ACADEMIC}，开题、发布节奏与授牌预热"),
    ], [4.5, 12])
    wbod(doc, "2.2 明确不承接、不承诺")
    for t in [
        "不取代同浦汇作为园区产业招服入口与客户关系主人",
        "不做租金去化对赌，也不做「必须租出去」的必要性条款（办公约3.3元/㎡/天高于周边）",
        "不承诺必带外资企业、不承诺领事必到场",
        "不把出海专题、领事到访、挂牌仪式装进30万元年包",
        "不向企业承诺高新/载体/补贴必然获批",
        "园区销售促成仍由园区负责；我方与同浦汇负责带到场、建群、回访摘要",
    ]:
        wbullet(doc, t)

    wh(doc, "三、角色分工")
    wtable(doc, [
        ("主体", "在本承接方案中的责任"),
        (PARTY_A_SHORT, "园区接口、在谈客户、带客名单协同、建群与回访台账、向园区转交促成线索"),
        (PARTY_B_SHORT, "30场策划与现场执行、服中心挂牌筹备、政策申报材料、月报验收"),
        ("科技企业联合会", "联合协办、行业组织协同、活动背书"),
        ("复旦住房政策研究中心", "学术支持、白皮书、专家与行业数据、海外专题研究（按场次确认）"),
        (f"{PARK} / 园区", "场地物业、导视、样板间、在谈名单、销售促成"),
        (OWNER, "空间业主；重大口径与场地条件确认"),
    ], [4.5, 12])
    wbod(doc, f"同浦汇联系人：{CONTACT_TPH}")
    wbod(doc, "空间业主：杨浦科创集团（杨浦国资）。审核相关方：中建相关工程/投资审核方（需过审，不写入合同业主栏）。")

    wh(doc, "四、工作包A：30场活动全案")
    wbod(doc, "4.1 锁版商务口径（对园区）")
    wtable(doc, [
        ("项目", "口径"),
        ("产品", "30场 / 30万元（固定执行费）"),
        ("周期", PERIOD),
        ("人数 / 负责人", "每场≤30人（可核验）；负责人占比≤30%，不做100%承兑"),
        ("触达", "全年约600+人次"),
        ("付款", "签约后7日内50%启动款；Q3、Q4、次年Q1、次年Q2各10%；年终收官10%尾款"),
        ("验收", "季度场次完成度 + 月报验收"),
        ("门票赞助", "可冲抵成本，不替代打包价"),
    ], [4, 12.5])
    wbod(doc, "4.2 新赛道下的六条线")
    wtable(doc, [
        ("线条", "场次", "时间窗", "作用"),
        ("启动 E", "5", "8月点火", "赛道切换、渠道、金融、政府、进博承接"),
        ("智能建造 A", "7", "全年能力线", "机器人、模块化、零碳、测绘、黑客松"),
        ("政策 B", "4", "申报入口", "高新、创新券、YOUNG立方、冲刺营"),
        ("出海准备 C", "6", "国内可交付", "案例、高校、华东华南、产业地图"),
        ("展示 D", "5", "空间产品化", "500㎡展厅、白皮书、研学、联展"),
        ("收官 F", "3", "可核验", "Demo Day、圆桌、年终验收"),
    ], [3.5, 2, 3.5, 7.5])
    wbod(doc, "4.3 逐场排期")
    ev_rows = [("编号", "月份", "线条", "活动名称", "本场作用")]
    for code, month, line, name, role in EVENTS:
        ev_rows.append((code, month, line, name, role))
    wtable(doc, ev_rows, [1.8, 2.2, 2.2, 5.5, 5])
    wbod(doc, "4.4 单场标准与转化闭环")
    for t in [
        "执行节奏：T-14策划邀约 → T-7确认名单 → T-1场地物料 → D日执行 → T+1~7建群/回访摘要",
        "转化闭环：到场 → 建群 → 看场 → 同浦汇回访 → 园区销售促成 → 复盘回流下一场",
        "交付物：方案、签到、现场执行、月报；回访只交摘要知会，不要求园区逐户回访",
    ]:
        wbullet(doc, t)
    wbod(doc, "4.5 不进年包的加购项")
    wtable(doc, [
        ("项目", "内容", "计价原则"),
        ("出海专题执行", "境外对接、订单撮合、海外考察", "10–20万/场次或专项，一事一议"),
        ("领事到访接待", "驻沪领馆匹配目标国", "按场次另计价，与挂牌分开"),
        ("挂牌仪式", "服中心 / 复旦战略合作授牌", "3万/个起，挂牌≠领事活动"),
        ("战略签约大会升级", "媒体、会场、嘉宾规格升级", "在预热沙龙之外另报价"),
    ], [4, 6, 6.5])

    wh(doc, "五、工作包B：科技企业服务中心筹备")
    wbod(doc,
         "90天内，把「上海市杨浦区科技企业服务中心」在创智汇做成可挂牌、可辅导、可申报的服务入口。"
         "挂牌仪式另计价；筹备过程不另向同浦汇收取费用。政策收益按第八章分成。")
    wtable(doc, [
        ("窗口", "主题", "交付"),
        ("D0–15", "叙事与主体", "统一新赛道一页纸；确认挂牌主体、2人配置、创智汇接待位/500㎡展区界面"),
        ("D16–45", "制度与申报", "服务制度、辅导SOP、载体认定材料包、创新券机构申请、高新路径清单"),
        ("D46–75", "产品化", "展陈大纲、零碳研学模型、白皮书开题、首批可辅导企业名单"),
        ("D76–90", "可过审包", "给审核方工程部门的初步方案、挂牌仪式报价（另计价）、Q4场次点火"),
    ], [3, 3.5, 10])
    wtable(doc, [
        ("空间", "口径"),
        ("项目", "创智汇（一期）｜上海·杨浦·五角场·创智天地｜约6600㎡"),
        ("3F", "约2850㎡　孵化器 + 办公 → 智能建造办公、装备体验、辅导接待"),
        ("5F", "约3670㎡　展厅 + 贸易 → 出海展陈、模块化样品、新材料、培训沙龙"),
        ("专题展区", "约500㎡　建筑产业出海主题展示（会上口径）"),
        ("租金 / 物业", "办公约3.3元/㎡/天｜物业约13.8元/㎡/月｜待租约8间（沟通口径）"),
    ], [3.5, 13])
    wbod(doc, "租金约3.3元高于周边，故本承接方案不做租赁对赌、不做租赁必要性。建议园区给予1–3个月免租，由同浦汇与园区线下谈，不写入服中心对同浦汇的承诺。")

    wh(doc, "六、给审核方的三句口径")
    for t in [
        "我们不是在招「付不起租金的施工队来租办公室」。",
        "我们是在用科技企业服务中心，把智能建造产品、模块化建筑和绿色低碳能力组织成可出海的集群。",
        "业主是杨浦科创集团；政府对接使用科技与产业服务口径，不以「中建」名义包装。",
    ]:
        wbullet(doc, t)

    wh(doc, "七、可核验指标")
    wtable(doc, [
        ("指标", "标准"),
        ("场次", "12个月内完成30场可核验活动"),
        ("人数", "单场签到≤30人，全年触达约600+"),
        ("月报", "每月1份，含名单摘要、照片、线索条数"),
        ("服中心", "90天内完成挂牌材料与制度包，具备申报条件"),
        ("转化", "每场输出看场线索，由同浦汇回访、园区促成；不设租金对赌"),
    ], [3.5, 13])

    wh(doc, "八、商务与内部结算")
    wbod(doc, "8.1 对园区：活动年包30万元；付款节点见第四章；不做租金对赌。招商佣金按2个月净租金（首年不重复），由园区销售闭环触发。")
    wbod(doc, "8.2 同浦汇 × 服中心（建议，供确认并写入协议）")
    wtable(doc, [
        ("收入类型", "建议分成", "说明"),
        ("活动年包30万", "服中心70% / 同浦汇30%", "活动年包对园区仍为30万；内部按执行与接口拆分"),
        ("政策申报收益", "同浦汇38% / 服中心62%", "以资金到账为准；两套分成互不混用"),
        ("服中心筹备", "不另向同浦汇收费", "向政策与园区要的是牌照与申报入口"),
        ("出海 / 领事 / 挂牌", "一事一议", "与年包、政策分成分开签署"),
    ], [4, 5, 7.5])

    wh(doc, "九、下一步（请同浦汇确认）")
    for t in [
        "1. 同浦汇确认：承接范围（30场 + 服中心筹备）与不承接边界无异议",
        "2. 确认内部结算：活动年包70/30、政策收益38/62",
        "3. 确认8–9月前8场档期与场地",
        "4. 确认向审核方提交的业主口径（杨浦科创集团）与赛道叙事",
        "5. 签约后7日内启动款到位，T+14交付细化执行手册",
    ]:
        wbullet(doc, t)

    wh(doc, "十、配套文件")
    for t in [
        f"PPT：《{PARTY_A_SHORT}_30场活动与科技企业服务中心筹备_业务承接策划案.pptx》",
        f"Excel：《{PARTY_A_SHORT}_30场活动与科技企业服务中心筹备_执行台账.xlsx》",
        f"协议：《{PARTY_A_SHORT}_创智汇30场活动与服中心筹备_业务承接合作协议.docx》（可直接签署）",
        "本文Word为可打印、可批注的完整策划案正文（复核定稿）。",
    ]:
        wbullet(doc, t)

    wp(doc, "（正文完）", align=WD_ALIGN_PARAGRAPH.CENTER, after=12)
    wp(doc, PARTY_B, align=WD_ALIGN_PARAGRAPH.RIGHT, after=2)
    wp(doc, DATE_CN, align=WD_ALIGN_PARAGRAPH.RIGHT, after=2)

    path = OUT / f"{PARTY_A_SHORT}_30场活动与科技企业服务中心筹备_业务承接策划案.docx"
    doc.save(path)
    return path


# —— PPT ——
BG = RGBColor(0x1A, 0x11, 0x40)
PANEL = RGBColor(0x27, 0x1A, 0x4D)
INK = RGBColor(0xF5, 0xF1, 0xFF)
MUT = RGBColor(0xC6, 0xBA, 0xEA)
SOFT = RGBColor(0x94, 0x86, 0xC4)
ACC = RGBColor(0x8B, 0x7B, 0xFF)
GOLD = RGBColor(0xE9, 0xC2, 0x7C)
ROW_A = RGBColor(0x2A, 0x1C, 0x52)
ROW_B = RGBColor(0x23, 0x17, 0x48)
SW, SH = Inches(13.333), Inches(7.5)
ML, CW = Inches(0.55), Inches(12.2)


def _rgb_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_text(slide, left, top, width, height, text, *, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = PPt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Microsoft YaHei"
    return box


def _card(slide, left, top, width, height, fill=PANEL):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    _rgb_fill(sh, fill)
    sh.adjustments[0] = 0.08
    return sh


def _new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    _rgb_fill(bg, BG)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.08))
    _rgb_fill(bar, ACC)
    return s


def _footer(s, page, total=16):
    _add_text(s, ML, Inches(7.05), CW, Inches(0.3),
              f"{PARTY_B_SHORT} × {PARTY_A_SHORT} · 业务承接策划案（复核定稿）· 2026.09　　{page:02d} / {total:02d}",
              size=10, color=SOFT)


def _header(s, no, title, sub=""):
    _add_text(s, ML, Inches(0.22), Inches(1.2), Inches(0.35), no, size=14, color=GOLD, bold=True)
    _add_text(s, Inches(1.5), Inches(0.18), Inches(10), Inches(0.4), title, size=22, color=INK, bold=True)
    if sub:
        _add_text(s, ML, Inches(0.62), CW, Inches(0.3), sub, size=12, color=MUT)


def _table(slide, left, top, width, height, rows, col_w=None, font_size=11):
    nrows, ncols = len(rows), len(rows[0])
    table_shape = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    table = table_shape.table
    if col_w:
        for i, w in enumerate(col_w):
            table.columns[i].width = w
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = PPt(font_size)
                p.font.name = "Microsoft YaHei"
                p.font.bold = ri == 0
                p.font.color.rgb = INK if ri == 0 else MUT
                p.alignment = PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = ROW_A if ri == 0 else (ROW_B if ri % 2 else PANEL)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return table_shape


def build_pptx() -> Path:
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    s = _new_slide(prs)
    _add_text(s, ML, Inches(1.3), CW, Inches(0.35), f"提交对象　{PARTY_A}", size=14, color=GOLD, bold=True)
    _add_text(s, ML, Inches(1.85), CW, Inches(0.9),
              f"{PARK}30场活动与科技企业服务中心筹备\n业务承接合作策划案（复核定稿）", size=28, color=INK, bold=True)
    _add_text(s, ML, Inches(3.0), CW, Inches(0.55),
              "承接范围：30场可核验活动　＋　科技企业服务中心挂牌筹备\n新赛道：智能建造 × 建筑产业出海　｜　政策包装：现代服务业 + 新一代信息技术",
              size=14, color=MUT)
    for i, (a, b) in enumerate([
        ("30场", "年度活动包"), ("30万", "对园区打包价"), ("≤30人", "单场核验"),
        ("≤30%", "负责人占比"), ("90天", "服中心筹备"), ("70/30", "执行/接口"),
    ]):
        x = ML + Inches(i * 2.05)
        _card(s, x, Inches(4.0), Inches(1.9), Inches(1.35))
        _add_text(s, x + Inches(0.1), Inches(4.15), Inches(1.7), Inches(0.55), a, size=22, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
        _add_text(s, x + Inches(0.1), Inches(4.75), Inches(1.7), Inches(0.4), b, size=12, color=MUT, align=PP_ALIGN.CENTER)
    _add_text(s, ML, Inches(5.7), CW, Inches(0.6),
              f"{PARTY_B}　｜　联合 {PARTY_C}\n学术支持：{ACADEMIC}　｜　{DATE_CN}　｜　周期 {PERIOD}", size=12, color=SOFT)
    _add_text(s, ML, Inches(6.5), CW, Inches(0.3),
              "机密 · 仅供同浦汇合作沟通，不作为对园区/投资方的承诺函", size=11, color=SOFT)
    _footer(s, 1)

    s = _new_slide(prs)
    _header(s, "01", "背景", "昨天沟通后，需要立刻承接的两件事　｜　2026-08-31 赛道调整交流")
    _card(s, ML, Inches(1.15), Inches(5.9), Inches(5.4))
    _card(s, Inches(6.8), Inches(1.15), Inches(5.9), Inches(5.4))
    _add_text(s, ML + Inches(0.25), Inches(1.3), Inches(5.4), Inches(0.35), "赛道为什么改", size=16, color=GOLD, bold=True)
    _add_text(s, ML + Inches(0.25), Inches(1.75), Inches(5.4), Inches(4.5),
              "1. 上周已定换赛道：由AI+IP转为建筑设计+智能建造，匹配同济设计与中建相关方主业。\n\n"
              "2. 按「现代服务业+新一代信息技术」包装，科技赋能建筑，不与高新/载体申报本质冲突。\n\n"
              "3. 目标不是招施工企业租办公室，而是做建筑产业出海集群。\n\n"
              "4. 传统施工回款弱于AI赛道，按招租逻辑难过审、也难转化。", size=13, color=MUT)
    _add_text(s, Inches(7.05), Inches(1.3), Inches(5.4), Inches(0.35), "新方向怎么落地", size=16, color=GOLD, bold=True)
    _add_text(s, Inches(7.05), Inches(1.75), Inches(5.4), Inches(4.5),
              "1. 出海优先城市化率较低的一带一路、东南亚、中亚、非洲、中东。\n\n"
              "2. 抓手：绿色低碳 + 机器人 + 模块化。\n\n"
              "3. 展厅约500㎡做建筑出海主题展示。\n\n"
              "4. 领事到访不纳入30场年包；合同口径已变、运营成本已发生。", size=13, color=MUT)
    _footer(s, 2)

    s = _new_slide(prs)
    _header(s, "02", "范围", "不是替换同浦汇的园区入口地位，而是把执行和筹备接下来")
    for i, (no, title, desc) in enumerate([
        ("01", "30场活动全案执行", "策划、档期、邀约协同、现场执行、月报、可核验交付"),
        ("02", "科技企业服务中心筹备", "挂牌路径、制度人员、载体认定、创新券、高新SOP、展陈研学"),
        ("03", "新赛道叙事与审核材料", "把昨日共识写成可过审一页纸 + 90天工作包"),
        ("04", "学术与白皮书协同", f"对接{ACADEMIC}，开题、发布与授牌预热"),
    ]):
        y = Inches(1.15) + Inches(i * 0.85)
        _card(s, ML, y, CW, Inches(0.75))
        _add_text(s, ML + Inches(0.2), y + Inches(0.1), Inches(0.7), Inches(0.5), no, size=18, color=GOLD, bold=True)
        _add_text(s, ML + Inches(1.0), y + Inches(0.08), Inches(10), Inches(0.3), title, size=15, color=INK, bold=True)
        _add_text(s, ML + Inches(1.0), y + Inches(0.38), Inches(10.5), Inches(0.3), desc, size=12, color=MUT)
    _footer(s, 3)

    s = _new_slide(prs)
    _header(s, "03", "角色", "同浦汇管关系，服中心管交付；园区管销售促成")
    for i, (name, desc) in enumerate([
        (PARTY_A_SHORT, "园区接口、在谈客户、带客名单协同、建群与回访台账、向园区转交促成线索"),
        (PARTY_B_SHORT, "30场策划与现场执行、服中心挂牌筹备、政策申报材料、月报验收"),
        ("联合会", "联合协办、行业组织协同、活动背书"),
        ("复旦住房中心", "学术支持、白皮书、专家与行业数据、海外专题（按场次确认）"),
        (f"{PARK}/园区", "场地物业、导视、样板间、在谈名单、销售促成"),
        (OWNER, "空间业主；重大口径与场地条件确认"),
    ]):
        col, row = i % 3, i // 3
        x = ML + Inches(col * 4.15)
        y = Inches(1.25) + Inches(row * 2.5)
        _card(s, x, y, Inches(3.95), Inches(2.2))
        _add_text(s, x + Inches(0.2), y + Inches(0.25), Inches(3.5), Inches(0.4), name, size=16, color=GOLD, bold=True)
        _add_text(s, x + Inches(0.2), y + Inches(0.8), Inches(3.5), Inches(1.2), desc, size=13, color=MUT)
    _footer(s, 4)

    s = _new_slide(prs)
    _header(s, "04", "产品", "对园区报价不变　｜　对同浦汇把执行责任写清楚")
    _card(s, ML, Inches(1.2), Inches(5.95), Inches(5.3))
    _card(s, Inches(6.85), Inches(1.2), Inches(5.95), Inches(5.3))
    _add_text(s, ML + Inches(0.25), Inches(1.4), Inches(5.4), Inches(0.4), "工作包A　30场活动全案", size=16, color=GOLD, bold=True)
    _add_text(s, ML + Inches(0.25), Inches(1.95), Inches(5.4), Inches(4.2),
              f"周期：{PERIOD}　·　约每月2–3场\n\n打包：30场 / 30万元（门票赞助可冲抵成本）\n\n"
              "核验：每场≤30人，负责人≤30%，全年约600+\n\n主线：智能建造能力 + 出海准备（国内可交付场）\n\n"
              "付款：签约7日内50% + 四季度各10% + 年终10%\n\n不含：出海专题、领事到访、挂牌仪式", size=13, color=MUT)
    _add_text(s, Inches(7.1), Inches(1.4), Inches(5.4), Inches(0.4), "工作包B　服中心筹备", size=16, color=GOLD, bold=True)
    _add_text(s, Inches(7.1), Inches(1.95), Inches(5.4), Inches(4.2),
              "90天完成挂牌材料、制度、2人配置\n\n载体认定 + 创新券服务机构两项前置\n\n"
              "高新辅导SOP（企业自行申报，不承诺获批）\n\n500㎡建筑出海展陈大纲 + 零碳研学模型\n\n"
              "联合复旦白皮书开题与授牌预热\n\n筹备不另向同浦汇收费\n\n政策收益：同浦汇38% / 服中心62%", size=13, color=MUT)
    _footer(s, 5)

    s = _new_slide(prs)
    _header(s, "05", "活动结构", "30场按六条线重排，服务新赛道")
    for i, (a, b, c, d) in enumerate([
        ("启动 E", "5场", "8月点火", "赛道切换、渠道、金融、政府、进博"),
        ("智能建造 A", "7场", "全年能力线", "机器人、模块化、零碳、测绘、黑客松"),
        ("政策 B", "4场", "申报入口", "高新、创新券、YOUNG立方、冲刺营"),
        ("出海准备 C", "6场", "国内可交付", "案例、高校、华东华南、产业地图"),
        ("展示 D", "5场", "空间产品化", "500㎡展厅、白皮书、研学、联展"),
        ("收官 F", "3场", "可核验", "Demo Day、圆桌、年终验收"),
    ]):
        col, row = i % 3, i // 3
        x = ML + Inches(col * 4.15)
        y = Inches(1.3) + Inches(row * 2.55)
        _card(s, x, y, Inches(3.95), Inches(2.3))
        _add_text(s, x + Inches(0.2), y + Inches(0.25), Inches(3.5), Inches(0.35), a, size=16, color=GOLD, bold=True)
        _add_text(s, x + Inches(0.2), y + Inches(0.7), Inches(3.5), Inches(0.3), f"{b}　·　{c}", size=13, color=ACC, bold=True)
        _add_text(s, x + Inches(0.2), y + Inches(1.2), Inches(3.5), Inches(0.8), d, size=13, color=MUT)
    _footer(s, 6)

    for half, page, title in [
        (EVENTS[:16], 7, "排期（上）· 2026.08—2026.12　共16场"),
        (EVENTS[16:], 8, "排期（下）· 2027.01—2027.07　共14场"),
    ]:
        s = _new_slide(prs)
        _header(s, f"0{page - 1}", title, "每场≤30人　｜　负责人≤30%")
        rows = [["编号", "月份", "线条", "活动名称", "本场作用"]]
        for code, month, line, name, role in half:
            rows.append([code, month, line, name, role])
        _table(s, ML, Inches(1.05), CW, Inches(5.7), rows,
               col_w=[Inches(0.9), Inches(1.2), Inches(1.3), Inches(4.5), Inches(4.3)], font_size=10)
        _footer(s, page)

    s = _new_slide(prs)
    _header(s, "08", "打法", "单场可复制，转化不压到同浦汇一个人身上")
    _add_text(s, ML, Inches(1.15), CW, Inches(0.35),
              "T-14 策划邀约 → T-7 确认名单 → T-1 场地物料 → D日执行 → T+1~7 建群 / 回访摘要",
              size=14, color=GOLD, bold=True)
    for i, (a, b) in enumerate([("≤30人", "单场签到核验"), ("≤30%", "负责人占比上限"), ("600+", "全年触达人次"), ("月报", "方案/签到/执行摘要")]):
        x = ML + Inches(i * 3.15)
        _card(s, x, Inches(1.7), Inches(3.0), Inches(1.3))
        _add_text(s, x + Inches(0.1), Inches(1.85), Inches(2.8), Inches(0.45), a, size=22, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
        _add_text(s, x + Inches(0.1), Inches(2.4), Inches(2.8), Inches(0.35), b, size=12, color=MUT, align=PP_ALIGN.CENTER)
    _card(s, ML, Inches(3.3), CW, Inches(3.2))
    _add_text(s, ML + Inches(0.3), Inches(3.5), Inches(11.5), Inches(0.35),
              "转化闭环（活动只对带到场负责，不对租金去化对赌）", size=15, color=GOLD, bold=True)
    _add_text(s, ML + Inches(0.3), Inches(4.0), Inches(11.5), Inches(0.5),
              "到场 → 建群 → 看场 → 同浦汇回访 → 园区促成 → 复盘回流", size=16, color=INK, bold=True)
    _add_text(s, ML + Inches(0.3), Inches(4.7), Inches(11.5), Inches(1.5),
              "另计价不进30万包：出海专题（10–20万/场或专项）· 领事到访（按场次）· 挂牌仪式（3万/个起）· 战略签约升级\n\n"
              "同浦汇在闭环中多承担客户跟踪与回访台账；服中心重心放在下一场策划与现场。园区拿回访摘要即可，不要求逐户回访。",
              size=13, color=MUT)
    _footer(s, 9)

    s = _new_slide(prs)
    _header(s, "09", "服中心", "90天把科技企业服务中心筹备成可挂牌、可申报")
    for i, (a, b, c) in enumerate([
        ("D0–15", "叙事与主体", "统一新赛道一页纸；确认挂牌主体、2人配置、接待位/500㎡展区"),
        ("D16–45", "制度与申报", "服务制度、辅导SOP、载体认定材料、创新券机构、高新路径"),
        ("D46–75", "产品化", "展陈大纲、零碳研学模型、白皮书开题、首批可辅导企业"),
        ("D76–90", "可过审包", "给审核方初步方案、挂牌仪式报价（另计价）、Q4点火"),
    ]):
        x = ML + Inches(i * 3.15)
        _card(s, x, Inches(1.25), Inches(3.0), Inches(3.3))
        _add_text(s, x + Inches(0.15), Inches(1.45), Inches(2.7), Inches(0.35), a, size=14, color=GOLD, bold=True)
        _add_text(s, x + Inches(0.15), Inches(1.9), Inches(2.7), Inches(0.35), b, size=15, color=INK, bold=True)
        _add_text(s, x + Inches(0.15), Inches(2.45), Inches(2.7), Inches(1.8), c, size=12, color=MUT)
    _card(s, ML, Inches(4.8), CW, Inches(1.7))
    _add_text(s, ML + Inches(0.25), Inches(5.0), Inches(11.7), Inches(1.3),
              "政策上限（测算，不是保底）：载体年上限110万｜活动YOUNG立方年上限200万｜十年三部分合计上限约6920万\n"
              "前置条件：①项目成为经认定载体/成果转化服务平台　②服中心作为创新券服务机构\n"
              "筹备不另向同浦汇收费；挂牌仪式另计价。", size=13, color=MUT)
    _footer(s, 10)

    s = _new_slide(prs)
    _header(s, "10", "空间", "约6600㎡仍是创智汇，产品换成新赛道　｜　办公约3.3元/㎡/天　｜　物业约13.8元/㎡/月")
    for i, (a, b, c) in enumerate([
        ("3F　约2850㎡", "孵化器 + 办公", "智能建造办公 / 装备体验 / 辅导接待 / OPC工位"),
        ("5F　约3670㎡", "展厅 + 贸易", "建筑出海展陈 / 模块化样品 / 新材料 / 培训沙龙"),
        ("专题展区　约500㎡", "建筑产业出海主题展示", "先做成可参观、可讲解、可研学的主题馆"),
    ]):
        y = Inches(1.2) + Inches(i * 1.15)
        _card(s, ML, y, CW, Inches(1.05))
        _add_text(s, ML + Inches(0.25), y + Inches(0.15), Inches(4), Inches(0.35), a, size=15, color=GOLD, bold=True)
        _add_text(s, ML + Inches(4.5), y + Inches(0.15), Inches(3), Inches(0.35), b, size=14, color=ACC, bold=True)
        _add_text(s, ML + Inches(0.25), y + Inches(0.55), Inches(11.5), Inches(0.35), c, size=13, color=MUT)
    _card(s, ML, Inches(4.8), CW, Inches(1.7))
    _add_text(s, ML + Inches(0.25), Inches(5.0), Inches(11.7), Inches(1.3),
              "与招商的关系：租金约3.3元高于周边，故本方案不做租赁对赌、不做租赁必要性。活动任务是把对的人带到场、建群、看场。\n"
              "待租约8间由园区销售闭环；佣金按2个月净租金（首年不重复）。建议园区给1–3个月免租，线下谈。\n"
              "零碳体验馆可做成面向中小学的研学产品，形成不完全依赖补贴的运营入口。", size=13, color=MUT)
    _footer(s, 11)

    s = _new_slide(prs)
    _header(s, "11", "政策", "新赛道仍然走科技企业服务，不另开一套「地产补贴」故事")
    for i, (a, b) in enumerate([
        ("对企业", "· 先改经营范围：智能施工/绿色节能/数字化设计\n· 企业自行申报高新，服中心提供辅导SOP\n· 不承诺获批，成功费以到账为准"),
        ("对载体", "· 争取认定为成果转化服务平台/孵化基地\n· 服中心申请成为创新券服务机构\n· YOUNG立方活动按投入50%、年封顶200万申报"),
        ("对品牌", "· 联合复旦发布智能建造/出海白皮书\n· 战略合作签约+官方授牌（仪式另计价）\n· 进博、领馆作为加购，不写进年包承诺"),
    ]):
        x = ML + Inches(i * 4.15)
        _card(s, x, Inches(1.2), Inches(3.95), Inches(3.4))
        _add_text(s, x + Inches(0.2), Inches(1.4), Inches(3.5), Inches(0.4), a, size=16, color=GOLD, bold=True)
        _add_text(s, x + Inches(0.2), Inches(2.0), Inches(3.5), Inches(2.3), b, size=13, color=MUT)
    _card(s, ML, Inches(4.9), CW, Inches(1.6))
    _add_text(s, ML + Inches(0.25), Inches(5.1), Inches(11.7), Inches(1.2),
              "必须写进给审核方材料的三句话\n"
              "1. 我们不是在招「付不起租金的施工队来租办公室」。\n"
              "2. 我们是在用科技企业服务中心，把智能建造产品、模块化建筑和绿色低碳能力组织成可出海的集群。\n"
              "3. 业主是杨浦科创集团；政府对接用科技与产业服务口径，不用「中建」名义。", size=12, color=MUT)
    _footer(s, 12)

    s = _new_slide(prs)
    _header(s, "12", "商务", "对园区的报价一口价，付款只保留一套　｜　方案二已取消")
    _table(s, ML, Inches(1.2), CW, Inches(5.2), [
        ["项目", "口径"],
        ["活动年包", "30场 / 30万元（固定执行费）"],
        ["付款节点", "签约后7日内50%启动款；Q3、Q4、次年Q1、次年Q2各10%；年终收官10%尾款"],
        ["验收挂钩", "季度场次完成度 + 月报验收"],
        ["门票 / 赞助", "可冲抵成本，不替代打包价"],
        ["招商佣金", "2个月净租金（首年不重复），由园区销售闭环触发"],
        ["服中心筹备", "不另向同浦汇收取；向政策与园区要的是牌照与申报入口"],
    ], col_w=[Inches(2.5), Inches(9.7)], font_size=14)
    _footer(s, 13)

    s = _new_slide(prs)
    _header(s, "13", "结算", "同浦汇保留接口价值，服中心按执行结算　｜　两套分成互不混用")
    _card(s, ML, Inches(1.25), Inches(5.95), Inches(3.6))
    _card(s, Inches(6.85), Inches(1.25), Inches(5.95), Inches(3.6))
    _add_text(s, ML + Inches(0.25), Inches(1.45), Inches(5.4), Inches(0.4), "活动年包内部结算（建议）", size=16, color=GOLD, bold=True)
    _add_text(s, ML + Inches(0.25), Inches(2.1), Inches(5.4), Inches(2.4),
              "服中心　70%\n策划、邀约协同、现场、月报、可核验交付\n\n同浦汇　30%\n园区接口、客户跟踪、回访台账、线索转交", size=15, color=MUT)
    _add_text(s, Inches(7.1), Inches(1.45), Inches(5.4), Inches(0.4), "政策收益分成（锁版）", size=16, color=GOLD, bold=True)
    _add_text(s, Inches(7.1), Inches(2.1), Inches(5.4), Inches(2.4),
              "同浦汇　38%　｜　服中心　62%\n\n以资金到账为准\n十年上限测算≠保证到账\n成功费不向企业承诺", size=15, color=MUT)
    _card(s, ML, Inches(5.15), CW, Inches(1.4))
    _add_text(s, ML + Inches(0.25), Inches(5.35), Inches(11.7), Inches(1.0),
              "给同浦汇的一句话：您继续做创智汇面前的招服入口和客户主人；30场能不能办完、服中心能不能挂上牌，由服中心按本方案承接交付。\n"
              "这样您可以把精力放在回访和园区销售协同上，也有一份能交给审核方的执行答卷。", size=13, color=MUT)
    _footer(s, 14)

    s = _new_slide(prs)
    _header(s, "14", "过审", "下周就能带去审核方工程部门的材料包　｜　业主口径：杨浦科创集团")
    for i, (no, title, desc) in enumerate([
        ("01", "一页纸定位", "智能建造×建筑产业出海，科技赋能，不是地产招租"),
        ("02", "30场总表", "本PPT第7–8页 + Excel执行台账"),
        ("03", "服中心90天", "挂牌、制度、载体、创新券四件套时间表"),
        ("04", "商务一页", "30万、付款50+4×10+10、不做租金对赌"),
        ("05", "边界清单", "出海/领事/挂牌另计价；不承诺外企"),
        ("06", "复旦协同函", "白皮书开题意向 + 授牌预热节奏（正式签约另走）"),
    ]):
        col, row = i % 3, i // 3
        x = ML + Inches(col * 4.15)
        y = Inches(1.3) + Inches(row * 2.5)
        _card(s, x, y, Inches(3.95), Inches(2.2))
        _add_text(s, x + Inches(0.2), y + Inches(0.25), Inches(3.5), Inches(0.35), no, size=14, color=GOLD, bold=True)
        _add_text(s, x + Inches(0.2), y + Inches(0.7), Inches(3.5), Inches(0.4), title, size=16, color=INK, bold=True)
        _add_text(s, x + Inches(0.2), y + Inches(1.25), Inches(3.5), Inches(0.7), desc, size=13, color=MUT)
    _footer(s, 15)

    s = _new_slide(prs)
    _header(s, "15", "请同浦汇确认", "五件事确认后，7日内可启动　｜　并签署业务承接合作协议")
    for i, t in enumerate([
        "1. 同浦汇确认：承接范围（30场 + 服中心筹备）与不承接边界无异议",
        "2. 确认内部结算：活动年包70/30、政策收益38/62",
        "3. 确认8–9月前8场档期与场地",
        "4. 确认向审核方提交的业主口径（杨浦科创集团）与赛道叙事",
        "5. 签约后7日内启动款到位，T+14交付细化执行手册",
    ]):
        y = Inches(1.25) + Inches(i * 0.75)
        _card(s, ML, y, CW, Inches(0.65))
        _add_text(s, ML + Inches(0.3), y + Inches(0.15), Inches(11.5), Inches(0.4), t, size=15, color=INK)
    _add_text(s, ML, Inches(5.3), CW, Inches(0.8),
              f"配套：执行台账.xlsx　｜　策划案.docx　｜　业务承接合作协议.docx（可直接签署）\n联系人：{CONTACT_TPH}　｜　提出方：{PARTY_B}",
              size=13, color=SOFT)
    _footer(s, 16)

    path = OUT / f"{PARTY_A_SHORT}_30场活动与科技企业服务中心筹备_业务承接策划案.pptx"
    prs.save(path)
    return path


def build_agreement() -> Path:
    doc = Document()
    for s in doc.sections:
        s.top_margin = Cm(2.54)
        s.bottom_margin = Cm(2.54)
        s.left_margin = Cm(2.8)
        s.right_margin = Cm(2.8)

    wp(doc, f"{PARK}30场活动执行及科技企业服务中心筹备", font=HEI, size=18, bold=True,
       align=WD_ALIGN_PARAGRAPH.CENTER, after=4)
    wp(doc, "业务承接合作协议", font=HEI, size=20, bold=True, align=WD_ALIGN_PARAGRAPH.CENTER, after=10)
    wp(doc, "合同编号：____________________", size=11, align=WD_ALIGN_PARAGRAPH.RIGHT, after=12)

    wbod(doc, "本协议由下列双方于____年____月____日在上海市杨浦区签署：")
    wp(doc, f"甲方（委托方）：{PARTY_A}（以下简称“甲方”或“同浦汇”）", bold=True, after=2)
    wbod(doc, "住所地：上海市杨浦区______________________________")
    wbod(doc, "统一社会信用代码：____________________")
    wbod(doc, "法定代表人/负责人：__________　　联系人：__________　　电话：__________")
    wp(doc, f"乙方（承接方）：{PARTY_B}（以下简称“乙方”或“服中心”）", bold=True, after=2, before=6)
    wbod(doc, "住所地：上海市杨浦区______________________________")
    wbod(doc, "统一社会信用代码/登记证号：____________________")
    wbod(doc, "法定代表人/负责人：__________　　联系人：__________　　电话：__________")
    wbod(doc,
         f"鉴于甲方负责{PARK}（以下简称“本项目”）产业招服接口与客户关系管理，乙方具备活动执行与科技企业服务能力，"
         "双方就甲方委托乙方承接本项目年度30场活动全案执行及科技企业服务中心挂牌筹备事宜，经协商一致，订立本协议。")

    wh(doc, "第一条　定义与项目概况")
    wbod(doc, f"1.1　本项目指位于上海市杨浦区五角场·创智天地片区的{PARK}产业空间（沟通口径合计约6600㎡，具体以业主确认面积清单为准）。")
    wbod(doc, f"1.2　业主方为{OWNER}（或其指定项目公司）。园区销售促成由园区负责，不因本协议转移至乙方。")
    wbod(doc, "1.3　本项目赛道定位为「智能建造 × 建筑产业出海」；政策包装口径为「现代服务业 + 新一代信息技术（科技赋能建筑）」。")
    wbod(doc, f"1.4　学术支持单位为{ACADEMIC}；联合协办单位为{PARTY_C}。上述单位非本协议签约主体，其具体协作另行确认。")
    wbod(doc, "1.5　本协议附件包括：《业务承接策划案》《执行台账（30场总表）》及双方后续书面确认的补充文件，与本协议具有同等效力。")

    wh(doc, "第二条　合作内容与工作包")
    wbod(doc, "2.1　工作包A（30场活动全案）：乙方负责策划、档期协同、邀约协同、现场执行、月报及可核验交付；甲方负责园区接口、带客名单协同、建群与回访台账、向园区转交促成线索。")
    wbod(doc, "2.2　工作包B（服中心筹备）：乙方在签约生效后90日内完成挂牌材料与制度包，具备申报条件；挂牌仪式费用另议，筹备过程不另向甲方收费。")
    wbod(doc, "2.3　30场活动名称、月份与线条以附件《执行台账》为准；经双方书面同意可微调档期，但全年可核验场次不少于30场。")
    wbod(doc, "2.4　另计价事项（不纳入本协议年包价款）：出海专题执行、领事到访接待、挂牌仪式、战略签约大会升级等，一事一议另行签署确认单或补充协议。")

    wh(doc, "第三条　执行标准与验收")
    wbod(doc, "3.1　单场人数：每场签到人数原则上不超过30人，以签到表核验。")
    wbod(doc, "3.2　负责人占比：单场企业负责人（法人/CEO/创始人/业务VP/采购决策人等）占比目标不超过30%，不做100%承兑。")
    wbod(doc, "3.3　全年触达：约600人次以上（按实际签到累计）。")
    wbod(doc, "3.4　执行节奏建议：T-14策划邀约 → T-7确认名单 → T-1场地物料 → D日执行 → T+1~7建群/回访摘要。")
    wbod(doc, "3.5　交付物：单场方案、签到表、现场执行记录、每月月报（含名单摘要、照片、线索条数）；回访由甲方主责，乙方配合输出摘要知会园区。")
    wbod(doc, "3.6　验收：按季度场次完成度及月报进行验收；年终完成30场核验及年终月报后支付尾款。")

    wh(doc, "第四条　双方权利义务")
    wbod(doc, "4.1　甲方义务：按约付款；及时提供在谈客户与带客协同信息；确认场地档期与园区接口；对重大对外口径（含业主名义）予以书面确认。")
    wbod(doc, "4.2　乙方义务：按约完成工作包A、B；保证交付材料真实可核验；未经甲方同意不以甲方或园区名义作出本协议约定以外的承诺。")
    wbod(doc, "4.3　转化闭环：到场→建群→看场→甲方回访→园区销售促成→复盘回流。乙方对带到场、建群及线索摘要负责；不对租金去化结果对赌。")
    wbod(doc, "4.4　联合协办与学术支持由乙方协调对接，具体出席与成果按场次确认，费用如超出年包范围另行商定。")

    wh(doc, "第五条　价款、付款与内部结算")
    wbod(doc, "5.1　对园区/项目侧活动年包价款为人民币叁拾万元整（¥300,000.00），对应工作包A之30场执行（门票、赞助可冲抵成本，不替代该打包价）。")
    wbod(doc, "5.2　付款节点（唯一方案）：")
    for t in [
        "（1）本协议签署后7日内，支付50%启动款（人民币15万元）；",
        "（2）2026年第三季度结束后，支付10%（人民币3万元）；",
        "（3）2026年第四季度结束后，支付10%（人民币3万元）；",
        "（4）2027年第一季度结束后，支付10%（人民币3万元）；",
        "（5）2027年第二季度结束后，支付10%（人民币3万元）；",
        "（6）年终30场核验及年终月报完成后15日内，支付10%尾款（人民币3万元）。",
    ]:
        wbullet(doc, t)
    wbod(doc, "5.3　各季度付款与该季度场次完成度、月报验收挂钩；因甲方原因导致延期的，相应节点顺延且不视为乙方违约。")
    wbod(doc, "5.4　甲方与乙方内部结算建议（签署时确认，两套分成互不混用）：")
    wbullet(doc, "（1）活动年包：乙方70%、甲方30%；")
    wbullet(doc, "（2）因本项目产生的政策申报类收益（以资金到账为准）：甲方38%、乙方62%；")
    wbullet(doc, "（3）工作包B筹备：不另向甲方收费。")
    wbod(doc, "5.5　招商佣金按成交租户2个月净租金（首年不重复）计取，由园区销售闭环触发，不纳入本协议第5.1条年包价款，具体以甲方与园区/业主相关约定为准。")
    wbod(doc, "5.6　乙方指定收款账户：开户名____________________　开户行____________________　账号____________________")
    wbod(doc, "甲方指定收款账户（用于内部分配等）：开户名____________________　开户行____________________　账号____________________")

    wh(doc, "第六条　边界与不承诺事项")
    wbod(doc, "6.1　乙方不取代甲方作为园区产业招服入口与客户关系主人。")
    wbod(doc, "6.2　双方均不做租金去化对赌，也不设定「必须租出去」的必要性条款。")
    wbod(doc, "6.3　不承诺必带外资企业、不承诺领事必到场；出海专题、领事到访、挂牌仪式不纳入30万元年包。")
    wbod(doc, "6.4　不向企业承诺高新技术企业、载体认定或补贴必然获批；辅导成功费（如有）以资金到账为准。")
    wbod(doc, "6.5　对外材料涉及业主口径时，统一表述为杨浦科创集团；不以「中建」名义作为合同业主对外包装。")

    wh(doc, "第七条　保密")
    wbod(doc, "7.1　一方因履行本协议知悉的对方商业秘密、客户信息、报价、未公开方案等，均属保密信息，未经书面同意不得向第三方披露或用于本协议以外目的。")
    wbod(doc, "7.2　保密义务在本协议终止后继续有效三年。")

    wh(doc, "第八条　期限、变更与解除")
    wbod(doc, f"8.1　本协议服务期限为{PERIOD}，自____年____月____日起至____年____月____日止（与付款周期衔接，具体以填写日期为准）。")
    wbod(doc, "8.2　变更须双方签署书面补充协议。")
    wbod(doc, "8.3　一方严重违约经书面催告后15日内未改正的，守约方有权解除本协议，并按已完成工作据实结算；已收款超出部分应退还，不足部分应补足。")

    wh(doc, "第九条　违约责任与争议解决")
    wbod(doc, "9.1　任何一方违反本协议，应赔偿因此给对方造成的直接损失。")
    wbod(doc, "9.2　本协议适用中华人民共和国法律。因本协议引起的争议，双方先协商；协商不成的，任一方向本项目所在地有管辖权的人民法院提起诉讼。")

    wh(doc, "第十条　其他")
    wbod(doc, "10.1　本协议构成双方就本事项的完整约定，取代此前口头或书面磋商（已单独生效且不冲突的保密协议除外）。")
    wbod(doc, "10.2　本协议一式肆份，甲乙双方各执贰份，具有同等法律效力。本协议自双方法定代表人/负责人签字（或签章）并加盖公章之日起生效。")
    wbod(doc, "10.3　未尽事宜，双方可另行签订补充协议；补充协议与本协议冲突的，以补充协议为准。")

    wp(doc, "", after=8)
    wp(doc, "（以下无正文，为签署页）", align=WD_ALIGN_PARAGRAPH.CENTER, after=16)
    wp(doc, "签署页", font=HEI, size=14, bold=True, align=WD_ALIGN_PARAGRAPH.CENTER, after=18)

    wp(doc, f"甲方：{PARTY_A}（盖章）", bold=True, after=8)
    wp(doc, "法定代表人/负责人（签字）：________________", after=6)
    wp(doc, "签署日期：____年____月____日", after=18)

    wp(doc, f"乙方：{PARTY_B}（盖章）", bold=True, after=8)
    wp(doc, "法定代表人/负责人（签字）：________________", after=6)
    wp(doc, "签署日期：____年____月____日", after=18)

    wp(doc, "附件清单（签署时勾选）：□ 业务承接策划案　□ 执行台账（30场总表）　□ 收益分配确认书　□ 其他：__________",
       size=11, after=6)

    path = OUT / f"{PARTY_A_SHORT}_创智汇30场活动与服中心筹备_业务承接合作协议.docx"
    doc.save(path)
    return path


def main():
    ART.mkdir(parents=True, exist_ok=True)
    paths = [build_xlsx(), build_docx_plan(), build_pptx(), build_agreement()]
    for p in paths:
        shutil.copy2(p, ART / p.name)
        print("OK", p.name, p.stat().st_size)
    print("DONE", len(paths))


if __name__ == "__main__":
    main()
