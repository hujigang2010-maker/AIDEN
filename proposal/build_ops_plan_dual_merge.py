# -*- coding: utf-8 -*-
"""创智汇年度活动运营方案 · 双会资源并入融合版

主结构沿用年度运营方案 PPT：
  01 背景与总体思路 → 02 日历与节奏 → 03 六大主题 → 04 执行标准
  → 05 转化·分工·收费·案例 → 06 结语 → 附录（定义核映射）
并入：WAIC/ChinaJoy 招商向资源 + 30场定义核 + 修订口径（合并同类项）
"""
import os, shutil
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

HERE = os.path.dirname(os.path.abspath(__file__))
ART = "/opt/cursor/artifacts"
PHOTO = os.path.join(HERE, "cases", "photos")
FONT = "Microsoft YaHei"

INK = RGBColor(0x1A, 0x1A, 0x2E)
MUT = RGBColor(0x4A, 0x4A, 0x5A)
SOFT = RGBColor(0x8A, 0x8A, 0x9A)
BLUE = RGBColor(0x1E, 0x4D, 0x8C)
GOLD = RGBColor(0xC4, 0x9A, 0x3C)
LINE = RGBColor(0xD8, 0xDE, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF7, 0xF8, 0xFA)
ACC = RGBColor(0xE8, 0xF0, 0xFA)
GREEN = RGBColor(0x1F, 0x7A, 0x5C)

SW, SH = Inches(13.333), Inches(7.5)
ML, CW = Inches(0.55), Inches(12.23)

EVENTS = [
    dict(code="P1", name="创智汇项目推广日①·中介与渠道专场", month="2026年8月上旬", fmt="项目推广日",
         people="≤30人", dm="约50%", floor="5F+3F", lease="中介渠道带看与项目推介",
         agenda="项目介绍→空间导览→政策要点→一对一洽谈→建群",
         guests="房产/产业中介、渠道、园区招商",
         define="渠道点火场", src="园区自建", res="待租8间看场清单；双会转介话术", layer="L0"),
    dict(code="P2", name="创智汇项目推广日②·五大行与金融机构专场", month="2026年8月中旬", fmt="项目推广日",
         people="≤30人", dm="约50%", floor="5F沙龙", lease="金融资源链接与客户导入",
         agenda="项目路演→金融产品对接→需求座谈→看场→建群",
         guests="五大行及合作金融机构对公/科创条线",
         define="金融转介场", src="园区自建+投研", res="WAIC金融科技议题摘要；CJ专业观众画像", layer="L0"),
    dict(code="P3", name="创智汇项目推广日③·投研机构专场", month="2026年8月中旬", fmt="项目推广日",
         people="≤30人", dm="约50%", floor="5F圆桌", lease="投研视角放大项目影响力",
         agenda="项目定位→产业赛道研判→圆桌→看场→建群",
         guests="券商研究所、产业研究院、智库",
         define="叙事定调场", src="园区自建", res="WAIC赛道结论；CJ「与AI同游」判断", layer="L0"),
    dict(code="P4", name="创智汇项目推广日④·政府部门与载体协同专场", month="2026年8月下旬", fmt="项目推广日",
         people="≤30人", dm="约50%", floor="3F+5F", lease="政策与载体协同落地",
         agenda="载体介绍→政策协同→云创基地资源→看场→建群",
         guests="区相关部门、载体、云创基地、服中心",
         define="政策协同场", src="云创+科企联/服中心", res="算力券/创新券口径；YOUNG立方政策包", layer="L0"),
    dict(code="E1", name="WAIC成果承接·创智汇开放交流日", month="2026年8月下旬", fmt="开放交流日",
         people="≤30人", dm="约50%", floor="3F+5F", lease="衔接WAIC议题与园区看场",
         agenda="WAIC议题速递→园区承接→双展台→洽谈→建群",
         guests="AI/内容企业、媒体、云创基地代表",
         define="WAIC余热承接", src="WAIC总表", res="A/B线索抽样；ACT-001/002/006议题切片", layer="L1"),
    dict(code="A2", name="Agent智能体搭建一日营", month="2026年9月上旬", fmt="实训营",
         people="≤30人", dm="约50%", floor="3F OPC", lease="Agent初创落位",
         agenda="签到→工作流→实操→Demo→看场", guests="Agent工程师、初创团队",
         define="3F入孵核：Agent→OPC", src="WAIC·Agent", res="Agent池抽样；ACT-001缩尺", layer="L2"),
    dict(code="C1", name="ChinaJoy主题·数字娱乐与IP授权对接会", month="2026年9月中旬", fmt="对接会",
         people="≤30人", dm="约50%", floor="5F展示中心", lease="CJ赛道企业/IP方入驻与展位",
         agenda="赛道解读→IP授权→一对一撮合→看场→建群",
         guests="游戏/动漫/潮玩IP方、渠道、内容公司",
         define="5F内容核：IP授权对接", src="CJ·游戏/IP", res="中小工作室+IP方；可触达263抽样", layer="L2"),
    dict(code="F1", name="出海专题活动（另议·另计价）", month="档期另议", fmt="专题（另议）",
         people="≤30人", dm="约50%", floor="另议", lease="出海服务（不在年度活动包内）",
         agenda="另案确定", guests="另案确定",
         define="另议核：出海不进年度包", src="CJ国际展团（另案）", res="BTOB外资线索仅作另案素材", layer="另计价"),
    dict(code="A3", name="多模态智能体工作坊", month="2026年10月上旬", fmt="工作坊",
         people="≤30人", dm="约50%", floor="3F培训", lease="多模态团队入驻",
         agenda="案例→实操→共创→评审→看场", guests="多模态应用团队、产品负责人",
         define="3F技术核：多模态看场", src="WAIC·多模态/AIGC", res="AIGC子类抽样；内容创意论坛切片", layer="L2"),
    dict(code="E4", name="ChinaJoy主题日·数字娱乐生态交流", month="2026年10月中旬", fmt="主题日",
         people="≤30人", dm="约50%", floor="5F", lease="CJ生态企业集中看场与洽谈",
         agenda="开幕→赛道分享→展洽→看场→建群", guests="数字娱乐企业、内容工作室、渠道",
         define="5F生态核：数字娱乐看场", src="CJ·游戏风云", res="Express/独立游戏中小团队池", layer="L2"),
    dict(code="B2", name="AI治理与可信智能体沙龙", month="2026年10月下旬", fmt="沙龙",
         people="≤30人", dm="约50%", floor="3F", lease="合规型企业信任导入",
         agenda="治理框架→合规清单→案例→问答→对接", guests="合规顾问、企业法务/负责人",
         define="信任核：治理合规", src="WAIC·治理论坛", res="治理议题摘要；合规顾问短名单", layer="L2"),
    dict(code="A4", name="火山引擎×算力Infra实务营", month="2026年11月上旬", fmt="厂商联训",
         people="≤30人", dm="约50%", floor="3F+云创基地", lease="高算力企业定向看场",
         agenda="政策包→代金券→案例→诊断→看场", guests="火山引擎、云创基地、企业技术负责人",
         define="算力核：云创×厂商联训", src="WAIC·算力", res="ACT-003缩尺；算力展商池抽样", layer="L2"),
    dict(code="B4", name="创新券·算力券·模型券实务沙龙", month="2026年11月中旬", fmt="实务沙龙",
         people="≤30人", dm="约50%", floor="3F", lease="用券企业聚集",
         agenda="三券规则→核销→案例→开户→入驻洽谈", guests="券平台、云厂商、企业负责人",
         define="券务核：用券聚集转化", src="云创+WAIC算力", res="三券核销路径；云厂商联系人抽样", layer="L2"),
    dict(code="L1", name="领事到访接待（另计价）", month="档期另议", fmt="外事接待（另计价）",
         people="另议", dm="—", floor="会客厅/5F", lease="领事到访（三部分费用之外）",
         agenda="另案确定", guests="领事及相关方（另案）",
         define="外事另计价核（到访≠挂牌）", src="另案", res="不占用年度30场资源池", layer="另计价"),
    dict(code="A5", name="具身智能空间交互体验日", month="2026年12月上旬", fmt="体验日",
         people="≤30人", dm="约50%", floor="5F+3F", lease="具身/机器人团队看场",
         agenda="演示→讲解→场景→踩点→洽谈", guests="具身团队、高校实验室",
         define="双会叠加核：具身体验", src="WAIC具身+CJ Vision Future", res="ACT-002缩尺；场景观察+CJ前沿科技", layer="L3"),
    dict(code="F4", name="创智汇AI年度Demo Day", month="2026年12月中旬", fmt="路演日",
         people="≤30人", dm="约50%", floor="5F主展", lease="集中签约洽谈与媒体背书",
         agenda="开幕→精选路演→颁奖→洽谈", guests="投资人、链主、媒体、精选项目",
         define="年中收口核：精选签约", src="双会精选项目", res="年内导入A/B意向精选路演", layer="L4"),
    dict(code="D1", name="AIGC微短剧制片特训", month="2027年1月上旬", fmt="特训营",
         people="≤30人", dm="约50%", floor="3F+5F", lease="厂牌/工作室入驻",
         agenda="政策→制片→脚本→路演→看场", guests="导演、厂牌、平台方",
         define="内容制片核", src="WAIC·AIGC", res="影视内容论坛切片；AIGC厂牌抽样", layer="L2"),
    dict(code="A6", name="AI营销Agent实战营", month="2027年1月中旬", fmt="实战营",
         people="≤30人", dm="约50%", floor="3F", lease="营销科技公司看场",
         agenda="策略→Agent→素材→复盘→转化", guests="投放操盘手、产品经理",
         define="营销科技核", src="WAIC·企业服务营销", res="企业服务/营销池抽样", layer="L2"),
    dict(code="B1", name="YOUNG立方×智能伙伴政策沙龙", month="2027年2月中旬", fmt="政策沙龙",
         people="≤30人", dm="约50%", floor="5F沙龙", lease="内容/AI企业导入",
         agenda="政策→礼包→画像→诊断→看场", guests="政策宣讲、服中心、企业负责人",
         define="政策获客核", src="服中心+双会客群", res="政策礼包+双会意向企业复邀", layer="L2"),
    dict(code="B3", name="高企认定冲刺（AI企业专场）", month="2027年2月下旬", fmt="辅导会",
         people="≤30人", dm="约50%", floor="3F", lease="待认定企业带政策入驻",
         agenda="条件→材料→初筛→套餐→激励", guests="辅导顾问、待认定企业负责人",
         define="带政策入驻核", src="服中心", res="双会AI应用类待认定抽样", layer="L2"),
    dict(code="A7", name="OPC超级个体黑客松（春）", month="2027年3月上旬", fmt="黑客松",
         people="≤30人", dm="约50%", floor="3F", lease="获奖团队优先谈单元",
         agenda="开题→开发→路演→礼包→看场", guests="评委、投资人、Builders",
         define="OPC核：获奖谈单元", src="WAIC青年+CJ创作者", res="高校Builder；CJ AGS线索抽样", layer="L2"),
    dict(code="C2", name="高校成果转化·AI for Science日", month="2027年3月下旬", fmt="对接日",
         people="≤30人", dm="约50%", floor="3F+沙龙", lease="成果公司/实验室落户",
         agenda="成果路演→场景→承接→看单元→洽谈",
         guests="复旦/同济技转、教授团队、住房政策研究中心学者",
         define="成果转化核", src="学术支持+WAIC科研", res="复旦住房中心学术支持位；科研论坛切片", layer="L2"),
    dict(code="F3", name="通往AGI季度圆桌", month="2027年3月下旬", fmt="闭门圆桌",
         people="≤30人", dm="约50%", floor="5F沙龙", lease="研究型/模型团队",
         agenda="议题→圆桌→围炉→问答→看场", guests="学者、模型企业、投资人",
         define="研究型客户核", src="WAIC·大模型论坛", res="ACT-001闭门形态；模型公司抽样", layer="L2"),
    dict(code="D2", name="ChinaJoy主题·创作者内容首发①", month="2027年4月中旬", fmt="发布会",
         people="≤30人", dm="约50%", floor="5F", lease="数字娱乐/内容品牌问询",
         agenda="揭幕→发布→快闪→专访→招商通道", guests="创作者、媒体、渠道、CJ内容方",
         define="创作者首发核", src="CJ·创作者展区", res="创作者/社团抽样；渠道玩法迁移", layer="L2"),
    dict(code="C3", name="ChinaJoy主题·游戏周边与潮玩供应链对接", month="2027年4月下旬", fmt="对接会",
         people="≤30人", dm="约50%", floor="5F", lease="CJ供应链/周边企业展位落位",
         agenda="赛道介绍→供需对接→展位参观→报价→建群", guests="游戏周边、潮玩供应链、渠道采购",
         define="供应链核：周边/潮玩展位", src="CJ·魔玩/谷子", res="潮玩谷子高优先级；作品类示例", layer="L2"),
    dict(code="C4", name="专精特新·AI应用培育路演", month="2027年5月中旬", fmt="路演",
         people="≤30人", dm="约50%", floor="5F", lease="成长型AI补位",
         agenda="政策→路演→点评→金融→看场", guests="顾问、银行、基金、企业",
         define="成长补位核", src="WAIC·行业AI", res="ACT-006缩尺；行业AI池抽样", layer="L2"),
    dict(code="D3", name="ChinaJoy主题·数字娱乐市集体验日", month="2027年5月下旬", fmt="体验日",
         people="≤30人", dm="约50%", floor="5F", lease="优质摊主/工作室升级固定展位",
         agenda="布展→体验→交流→转正洽谈→建群", guests="内容工作室、周边品牌、达人",
         define="市集转化核：摊主升级展位", src="CJ·创作者/市集", res="小微品牌；市集→固定展位路径", layer="L2"),
    dict(code="C5", name="ChinaJoy主题·IP联名与衍生品撮合会", month="2027年6月中旬", fmt="撮合会",
         people="≤30人", dm="约50%", floor="5F", lease="IP联名/衍生品团队入驻",
         agenda="IP路演→联名模式→一对一→MOU→看场", guests="IP方、衍生品商、渠道",
         define="IP商业化核", src="CJ·IP/联名", res="老字号×IP联名案例迁移；衍生品抽样", layer="L2"),
    dict(code="D4", name="ChinaJoy主题·沉浸式数字娱乐联展", month="2027年6月下旬", fmt="联展（短展期）",
         people="≤30人/场次", dm="约50%", floor="5F", lease="闭幕洽谈集中转化",
         agenda="布展→预约观展→交流→闭幕洽谈", guests="联合内容方、渠道、品牌",
         define="视觉落地核：短展期转化", src="CJ·氛围落地", res="联合内容方短名单；预约制控流", layer="L2"),
    dict(code="D5", name="创作者首发②·衔接WAIC2027与ChinaJoy", month="2027年7月中旬", fmt="发布会",
         people="≤30人", dm="约50%", floor="5F", lease="旺季补位，双大会预热",
         agenda="发布→渠道→看场→年框→双会预热", guests="IP方、渠道、媒体",
         define="双会预热核", src="双会预热", res="全年双会线索复盘；预热下届双会", layer="L4"),
]
assert len(EVENTS) == 30
EMAP = {e["code"]: e for e in EVENTS}

THEMES = [
    dict(key="A", sec="03 · 主题 A", title="智能体与算力训练（3F基本盘）",
         blurb=["训练营负责聚人，政策/Demo Day负责转化", "资源：WAIC Agent/多模态/算力/具身池抽样", "重点场：A2 Agent一日营 → OPC看场"],
         codes=["A2", "A3", "A4", "A5", "A6", "A7"]),
    dict(key="B", sec="03 · 主题 B", title="政策与治理沙龙（进园理由）",
         blurb=["把补贴、资质、合规做成面对面小场", "紧接算力营做券务，年底决策窗口转化", "服中心现场可收件启动申报"],
         codes=["B2", "B4", "B1", "B3"]),
    dict(key="C", sec="03 · 主题 C", title="产业与ChinaJoy对接（合并同类项）",
         blurb=["原东莞潮玩/汕头玩具/扬州毛绒 → 统一为CJ主题", "一头高校成果，一头CJ中小工作室与供应链", "重点：C1 IP授权对接 / C2成果转化日"],
         codes=["C1", "C2", "C3", "C4", "C5"]),
    dict(key="D", sec="03 · 主题 D", title="AI内容与ChinaJoy人气（5F引擎）",
         blurb=["特训供给 → 首发交易 → 市集/联展客流", "与CJ创作者展区、WAIC AIGC错峰协同", "重点：D4短展期联展闭幕洽谈"],
         codes=["D1", "D2", "D3", "D4", "D5"]),
    dict(key="E", sec="03 · 主题 E", title="启动月：项目推广 + 双会承接",
         blurb=["8月四场项目推广日（中介/金融/投研/政府）", "E1 WAIC成果开放交流日（≤30人精准）", "E4 CJ主题日承接数字娱乐生态"],
         codes=["P1", "P2", "P3", "P4", "E1", "E4"]),
    dict(key="F", sec="03 · 主题 F", title="收官 · 圆桌 · 另计价事项",
         blurb=["F4 Demo Day：精选路演集中签约", "F3 AGI圆桌：研究型客户深谈", "F1出海 / L1领事：另议另计价，不进年度包"],
         codes=["F4", "F3", "F1", "L1"]),
]

CASES = [
    {
        "sec": "05 · 往期案例（一）训练营 / 沙龙类",
        "name": "杨「数」浦数字沙龙第七期：AI如何重塑企业DNA",
        "meta": "2025-06-30 · 美团上海综合指挥中心 · 政策/治理沙龙类",
        "orgs": "杨浦区委网信办；赛博院等",
        "companies": "美团；大众点评大模型团队等",
        "agenda": [("参观", "美团综合指挥中心数字化场景参观"), ("主题", "人机共生时代：AI如何重塑企业DNA"),
                   ("案例", "大众点评AI技术应用创新与落地"), ("政策", "《人工智能生成合成内容标识办法》等解读"),
                   ("互动", "答疑与企业一对一交流")],
        "photos": ["case_c987bd91.jpg", "case_59157ea5.jpg", "case_6e51af9e.jpg"],
        "map": "对应：B类政策沙龙、A类训练营",
    },
    {
        "sec": "05 · 往期案例（二）开放日 / 路演大场类",
        "name": "「融见科创·智启未来」人工智能专场路演暨投融资对接会",
        "meta": "2025-10 · 杨浦 · 邮储银行联合主办 · 路演对接类",
        "orgs": "杨浦科创促进会、邮储银行等",
        "companies": "复楚智能、卡房信息、一造科技、中科趋势、万笔千墨等",
        "agenda": [("开场", "促进会会长夏立城致辞与生态介绍"), ("主办", "邮储银行科技金融致辞"),
                   ("分享", "中邮证券AI行业趋势主题分享"), ("路演", "6+6+6+1：项目+投资人+金融机构+政府"),
                   ("收口", "区投促点评与招商邀请")],
        "photos": ["case_db1d5cac.jpg", "case_65931516.jpg", "case_cb0a3812.jpg"],
        "map": "对应：F4 Demo Day、项目推广日、路演类",
    },
]


def rect(s, x, y, w, h, fill=None, line=None, lw=0.8):
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    b.shadow.inherit = False
    if fill is None: b.fill.background()
    else: b.fill.solid(); b.fill.fore_color.rgb = fill
    if line is None: b.line.fill.background()
    else: b.line.color.rgb = line; b.line.width = Pt(lw)
    return b


def text(s, x, y, w, h, content, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = content; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = FONT
    return tb


def bullets(s, x, y, w, h, items, size=12, color=MUT, mark=BLUE):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4); p.line_spacing = 1.15
        r0 = p.add_run(); r0.text = "▪ "; r0.font.size = Pt(size); r0.font.color.rgb = mark; r0.font.name = FONT
        r = p.add_run(); r.text = it; r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = FONT
    return tb


def slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, SW, SH, fill=WHITE)
    rect(s, 0, 0, SW, Inches(0.06), fill=BLUE)
    return s


def header(s, sec, title, sub=None):
    text(s, ML, Inches(0.22), Inches(12), Inches(0.28), sec, size=12, color=BLUE, bold=True)
    text(s, ML, Inches(0.48), Inches(12), Inches(0.42), title, size=22, color=INK, bold=True)
    if sub:
        text(s, ML, Inches(0.92), Inches(12), Inches(0.28), sub, size=12, color=MUT)


def footer(s, page, total):
    text(s, ML, Inches(7.15), Inches(9.5), Inches(0.25),
         "上海创智汇 × 同浦汇 · 年度活动运营方案 · 双会资源并入融合版", size=10, color=SOFT)
    text(s, Inches(11.2), Inches(7.15), Inches(1.6), Inches(0.25), f"{page}/{total}", size=10, color=SOFT, align=PP_ALIGN.RIGHT)


def card(s, x, y, w, h, title, items, accent=BLUE):
    rect(s, x, y, w, h, fill=BG, line=LINE)
    rect(s, x, y, Inches(0.08), h, fill=accent)
    text(s, Emu(x + Inches(0.22)), Emu(y + Inches(0.12)), Emu(w - Inches(0.35)), Inches(0.32), title, size=14, color=INK, bold=True)
    bullets(s, Emu(x + Inches(0.22)), Emu(y + Inches(0.48)), Emu(w - Inches(0.4)), Emu(h - Inches(0.55)), items, size=12, mark=accent)


def table(s, y0, rows, col_w, sizes=None, rh=0.32):
    yy = y0
    if sizes is None: sizes = [10] * len(rows[0])
    for ri, row in enumerate(rows):
        xx = ML; h = Inches(0.38 if ri == 0 else rh)
        for ci, val in enumerate(row):
            cw = Emu(int(CW * col_w[ci]))
            bg = BLUE if ri == 0 else (ACC if ri % 2 == 0 else WHITE)
            fg = WHITE if ri == 0 else INK
            rect(s, xx, yy, cw, h, fill=bg, line=LINE, lw=0.5)
            text(s, Emu(xx + Inches(0.04)), Emu(yy + Inches(0.04)), Emu(cw - Inches(0.06)), Emu(h - Inches(0.06)),
                 str(val), size=sizes[ci], color=fg, bold=(ri == 0), anchor=MSO_ANCHOR.MIDDLE)
            xx = Emu(xx + cw)
        yy = Emu(yy + h)
    return yy


def add_pic(s, path, x, y, w, h):
    if os.path.exists(path):
        s.shapes.add_picture(path, x, y, width=w, height=h)
    else:
        rect(s, x, y, w, h, fill=ACC, line=LINE)


def mon(s):
    return s.replace("2026年", "").replace("2027年", "次")


def build():
    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
    TOTAL = 29
    n = 0

    def pg(s):
        nonlocal n
        n += 1
        footer(s, n, TOTAL)
        return n

    # 1 cover
    s = slide(prs)
    text(s, ML, Inches(1.35), Inches(12), Inches(0.3), "上海创智汇 × 同浦汇｜呈报园区", size=14, color=BLUE, bold=True)
    text(s, ML, Inches(1.8), Inches(12), Inches(0.55), "创智汇年度活动运营方案", size=34, color=INK, bold=True)
    text(s, ML, Inches(2.5), Inches(12), Inches(0.35), "30场 · 定义核与双会资源导入 · 2026.08—2027.07", size=16, color=GOLD, bold=True)
    text(s, ML, Inches(3.0), Inches(12), Inches(0.35),
         "AI+数字内容无界共创港　｜　3F智能体与算力组团 · 5F内容与IP组团", size=13, color=MUT)
    card(s, ML, Inches(3.6), Inches(6.0), Inches(3.05), "本册并入要点（合并同类项）", [
        "结构：沿用年度运营方案章节",
        "内容：并入WAIC+ChinaJoy招商向资源",
        "口径：每场≤30人 · 负责人约50%",
        "分工：同浦汇带客 / 园区销售",
        "主体：云创基地（国家级孵化器）",
    ], BLUE)
    card(s, Inches(7.2), Inches(3.6), Inches(5.55), Inches(3.05), "边界", [
        "只取双会招商贴近切片，不堆全量名单",
        "出海另议另计价；领事到访/挂牌另计价",
        "不做租金租赁对赌与必要性要求",
        "不承诺必带外资企业",
        "学术支持：复旦住房政策研究中心",
    ], GREEN)
    pg(s)

    # 2 toc
    s = slide(prs); header(s, "目录", "本方案结构（以年度运营方案PPT为主）")
    table(s, Inches(1.35), [
        ["章节", "内容", "说明"],
        ["01 背景与总体思路", "为什么做、双会如何导入、做成什么样", "含层次叠加"],
        ["02 年度日历与运营节奏", "30场排期总表 + 四阶段重心 + 双会摘要", "合并资源摘要进正文"],
        ["03 六大主题活动详述", "A训练·B政策·C产业/CJ·D内容·E启动·F收官", "潮玩类合并为CJ主题"],
        ["04 单场执行标准", "筹备→看场→建群→回访同一套打法", "带客/销售分工"],
        ["05 转化·分工·收费·案例", "转化指标、付款、往期案例", "取消对赌/核验表"],
        ["06 结语", "交付与园区协同", "—"],
        ["附录", "30场定义核与资源导入映射", "不重复堆名单"],
    ], [0.26, 0.46, 0.28], sizes=[12, 12, 12], rh=0.48)
    pg(s)

    # 3 background
    s = slide(prs)
    header(s, "01 · 背景与判断", "大会之后，热度需要承接地——双会资源导入园区一年",
           "主题响应：上海创智汇 · AI+数字内容无界共创港")
    card(s, ML, Inches(1.35), Inches(5.9), Inches(5.4), "判断", [
        "WAIC闭幕后：议题/嘉宾/线索需落到看场与洽谈",
        "ChinaJoy「与AI同游」：中小工作室/创作者/潮玩/具身可导入5F",
        "3F缺持续入孵节奏；5F缺持续内容人气",
        "因此：双会是资源库，30场是定义核",
        "每月两到三场，精而准，不断档",
    ], BLUE)
    card(s, Inches(7.2), Inches(1.35), Inches(5.55), Inches(5.4), "运营目标（可执行可核验）", [
        "30场：全年活动总量",
        "12个月：2026.08—2027.07",
        "每场≤30人；全年约600人次+",
        "负责人占比适度承兑约50%",
        "不承诺必带外资企业",
    ], GOLD)
    pg(s)

    # 4 approach
    s = slide(prs); header(s, "01 · 总体思路", "一条主线：把双会的热度，变成园区的一年")
    card(s, ML, Inches(1.35), Inches(5.9), Inches(5.4), "园区得到什么", [
        "议题变日常：Agent/算力/治理/内容每月有场",
        "空间有主题：3F吃WAIC能力线，5F吃CJ内容线",
        "具身/AIGC为双会叠加带，跨楼层联动",
        "考核有支撑：签到、建群、回访摘要可归档",
        "云创基地国家级孵化器背书可感知",
    ], BLUE)
    card(s, Inches(7.2), Inches(1.35), Inches(5.55), Inches(5.4), "招商得到什么", [
        "客群更准：按赛道抽样邀约，不做全量名单轰炸",
        "动作标准化：每场固定看场+建群",
        "我们带客，园区负责销售促成",
        "政策/券务在活动中打包给出，降低决策门槛",
        "租金高于周边均价，故不做租赁对赌",
    ], GREEN)
    pg(s)

    # 5 framework
    s = slide(prs)
    header(s, "01 · 总体框架", "六大主题 + 四层叠加（合并同类项）", "主题管内容；层次管资源导入顺序")
    table(s, Inches(1.25), [
        ["主题", "名称", "场次", "在全年中的角色"],
        ["A", "智能体与算力训练", "6场", "3F基本盘：营/坊/联训/体验/黑客松"],
        ["B", "政策与治理沙龙", "4场", "政策/券务/资质做成进园理由"],
        ["C", "产业与ChinaJoy对接", "5场", "高校成果+CJ工作室/供应链/IP（合并潮玩类）"],
        ["D", "AI内容与ChinaJoy人气", "5场", "5F人气：特训/首发/市集/联展"],
        ["E", "启动月·推广与双会承接", "6场", "8月项目推广日+WAIC承接+CJ主题日"],
        ["F", "收官·圆桌·另计价", "2场+另计价", "Demo Day收口；出海/领事另案"],
    ], [0.1, 0.3, 0.14, 0.46], sizes=[12, 12, 11, 12], rh=0.4)
    text(s, ML, Inches(5.0), CW, Inches(0.25), "资源叠加", size=13, color=BLUE, bold=True)
    table(s, Inches(5.25), [
        ["层次", "作用", "映射"],
        ["L0点火", "中介/金融/投研/政府", "P1–P4"],
        ["L1承接", "WAIC余热→看场", "E1"],
        ["L2主线", "3F能力 / 5F内容分轨", "A/B/C/D常规场"],
        ["L3叠加 / L4收口", "具身叠加；签约+下届预热", "A5 · F4/D5"],
    ], [0.16, 0.36, 0.48], sizes=[11, 11, 11], rh=0.3)
    pg(s)

    # 6-7 tables
    for subset, title, sub in (
        (EVENTS[:15], "30场总表（上）", "2026年8月—12月（含另计价占位）"),
        (EVENTS[15:], "30场总表（下）", "2027年1月—7月"),
    ):
        s = slide(prs); header(s, "02 · 年度日历", title, sub + "　｜　每场≤30人 · 负责人约50%")
        rows = [["编号", "活动名称", "档期", "形式", "人数", "层次", "招商落点"]]
        for e in subset:
            rows.append([e["code"], e["name"][:16], mon(e["month"])[:8], e["fmt"][:8],
                         e["people"][:6], e["layer"][:6], e["lease"][:12]])
        table(s, Inches(1.25), rows, [0.07, 0.28, 0.12, 0.12, 0.1, 0.1, 0.21], sizes=[9, 10, 9, 9, 9, 9, 9], rh=0.32)
        pg(s)

    # 8 rhythm
    s = slide(prs); header(s, "02 · 运营节奏", "四个阶段，各有重心")
    card(s, ML, Inches(1.35), Inches(5.9), Inches(2.5), "启动期｜8—10月", [
        "8月四场项目推广日点火", "E1承接WAIC余热", "Agent营 + CJ IP对接开局", "CJ主题日与治理沙龙立住基调",
    ], BLUE)
    card(s, Inches(7.2), Inches(1.35), Inches(5.55), Inches(2.5), "攻坚期｜11—12月", [
        "算力联训 + 三券沙龙", "具身双会叠加体验", "Demo Day精选签约收口", "领事另计价另案（如需）",
    ], GOLD)
    card(s, ML, Inches(4.1), Inches(5.9), Inches(2.5), "深耕期｜1—3月", [
        "AIGC特训 / 营销Agent", "政策沙龙 + 高企冲刺", "黑客松 + 高校成果日", "AGI圆桌深谈研究型客户",
    ], GREEN)
    card(s, Inches(7.2), Inches(4.1), Inches(5.55), Inches(2.5), "收获与预热｜4—7月", [
        "CJ创作者首发 / 供应链对接", "市集转正 / IP联名 / 联展", "专精特新AI路演补位", "D5预热WAIC2027与下届CJ",
    ], BLUE)
    pg(s)

    # 9 dual resource
    s = slide(prs)
    header(s, "02 · 双会资源导入（摘要）", "只取招商贴近可调用规模与优先策略",
           "来源：WAIC2026全量资源整合总表 · ChinaJoy2026招商引资总表")
    card(s, ML, Inches(1.3), Inches(5.9), Inches(5.45), "WAIC → 喂养 3F", [
        "参展商963 · 品牌库4262 · 联系人3463 · A/B线索500",
        "主赛道：算力芯片 / 具身 / Agent / AIGC / 行业AI",
        "ACT缩尺：闭门会·具身体验·算力峰会·行业AI圆桌",
        "用法：线索抽样邀约；论坛议题→≤30人沙龙",
        "不做：全量展商粘贴进PPT",
    ], BLUE)
    card(s, Inches(7.2), Inches(1.3), Inches(5.55), Inches(5.45), "ChinaJoy → 喂养 5F", [
        "主题「与AI同游」· 可触达优先263 · 精编422",
        "最高：中小游戏工作室（Express/Next Play/AGS）",
        "高：Vision Future具身 · 潮玩谷子 · 创作者经济",
        "中腰部选址灵活；不硬追大厂总部搬迁",
        "原潮玩产业集群场次合并为CJ主题内容",
    ], GOLD)
    pg(s)

    # 10-15 themes
    for t in THEMES:
        s = slide(prs); header(s, t["sec"], t["title"])
        card(s, ML, Inches(1.3), Inches(12.23), Inches(1.55), "主题说明（合并后）", t["blurb"], BLUE if t["key"] != "F" else GOLD)
        rows = [["编号", "活动", "档期", "形式", "招商落点", "资源导入"]]
        for c in t["codes"]:
            e = EMAP[c]
            rows.append([e["code"], e["name"][:14], mon(e["month"])[:8], e["fmt"][:8], e["lease"][:12], e["res"][:18]])
        table(s, Inches(3.05), rows, [0.07, 0.24, 0.12, 0.12, 0.2, 0.25], sizes=[9, 10, 9, 9, 9, 9], rh=0.38)
        pg(s)

    # 16 execution
    s = slide(prs); header(s, "04 · 单场执行标准", "每场活动，同一套打法")
    table(s, Inches(1.35), [
        ["节点", "关键动作", "责任方", "交付物"],
        ["T-14", "锁定嘉宾议程物料 + 本场资源切片", "同浦汇", "议程/嘉宾确认/报名"],
        ["T-7", "邀约确认 + 场地报备", "同浦汇；园区场地确认", "名单/场地单"],
        ["T-1", "彩排 + 看场动线", "同浦汇；园区样板间", "彩排记录"],
        ["D日", "签到→主区→看场→建群", "同浦汇执行；园区支持", "签到/意向/群"],
        ["T+1~7", "带客回访", "同浦汇主责", "回访摘要知会园区"],
        ["销售促成", "房源条件与成交", "园区负责销售部分", "合同/交割"],
    ], [0.12, 0.34, 0.3, 0.24], sizes=[11, 12, 11, 11], rh=0.55)
    pg(s)

    # 17 conversion + division (merged)
    s = slide(prs); header(s, "05 · 转化与分工", "到场→建群→看场→销售促成（合并同类项）")
    card(s, ML, Inches(1.35), Inches(5.9), Inches(5.4), "转化动作", [
        "D日：签到分层并建群",
        "T+1：看场预约与资料入群",
        "T+1~7：同浦汇持续带客回访",
        "园区：群内发布房源/政策/条件，负责销售促成",
        "月度核对意向与成交进展",
    ], GREEN)
    card(s, Inches(7.2), Inches(1.35), Inches(5.55), Inches(5.4), "责任边界", [
        "同浦汇：策划、邀约、带客、建群、回访摘要",
        "园区：场地物业、样板间、销售成交、合同交割",
        "云创基地：主体运营范围与算力协同",
        "科企联/服中心：活动载体支持与政策手续",
        "复旦住房中心：学术支持",
    ], BLUE)
    pg(s)

    # 18 deliverables + fee (merged)
    s = slide(prs); header(s, "05 · 交付与收费", "交给园区的东西 · 付款唯一方案")
    table(s, Inches(1.25), [
        ["周期", "交付物"],
        ["签约后2周", "年度方案细化版 + 排期日历 + 资源导入台账模板"],
        ["每场后7天", "执行卡/议程/签到/意向/动线/照片通稿 + 回访摘要 + 建群清单"],
        ["每月", "活动数据月报 + 下月排期确认单"],
        ["每季/年末", "复盘、台账、双会资源消耗与补库建议"],
    ], [0.2, 0.8], sizes=[12, 12], rh=0.4)
    text(s, ML, Inches(4.0), CW, Inches(0.28), "收费与付款（合并后唯一方案）", size=13, color=BLUE, bold=True)
    table(s, Inches(4.3), [
        ["项目", "口径"],
        ["活动策划执行服务费", "年度打包（金额线下确认）；含常规30场"],
        ["租金/物业", "3.3元/㎡/天；物业13.8元/㎡/月；建议免租期1–3个月（面议）"],
        ["租赁对赌", "不做（租金高于周边均价）；不做租赁必要性要求"],
        ["出海/领事", "另议另计价；挂牌与到访分属不同性质"],
        ["付款", "签约后7日50%；Q3/Q4/次年Q1/次年Q2各10%；年终10%"],
    ], [0.28, 0.72], sizes=[11, 11], rh=0.34)
    pg(s)

    # 19-20 cases
    for c in CASES:
        s = slide(prs); header(s, c["sec"], c["name"], c["meta"])
        for j, fn in enumerate(c["photos"][:3]):
            add_pic(s, os.path.join(PHOTO, fn), Emu(ML + j * Inches(4.1)), Inches(1.35), Inches(3.95), Inches(2.25))
        rect(s, ML, Inches(3.8), Inches(7.5), Inches(2.95), fill=BG, line=LINE)
        text(s, Emu(ML + Inches(0.2)), Inches(3.9), Inches(7.1), Inches(0.3), "当天议程", size=13, color=BLUE, bold=True)
        yy = Inches(4.3)
        for t0, d0 in c["agenda"]:
            text(s, Emu(ML + Inches(0.2)), yy, Inches(1.0), Inches(0.28), t0, size=11, color=GOLD, bold=True)
            text(s, Emu(ML + Inches(1.2)), yy, Inches(6.0), Inches(0.28), d0, size=11, color=INK)
            yy = Emu(yy + Inches(0.28))
        rect(s, Inches(8.3), Inches(3.8), Inches(4.5), Inches(2.95), fill=ACC, line=LINE)
        text(s, Inches(8.45), Inches(3.9), Inches(4.2), Inches(0.3), "可见机构/企业", size=13, color=BLUE, bold=True)
        text(s, Inches(8.45), Inches(4.3), Inches(4.2), Inches(2.2),
             f"{c['orgs']}\n\n企业侧：{c['companies']}\n\n映射：{c['map']}", size=12, color=INK)
        pg(s)

    # 21 closing
    s = slide(prs); header(s, "06 · 结语", "把一年做满，双会热度落在园区")
    card(s, ML, Inches(1.35), Inches(5.9), Inches(5.4), "我们将交付", [
        "30场按日历执行（每场≤30人）",
        "8月项目推广日四场点火 + WAIC承接",
        "ChinaJoy主题内容场次（合并原潮玩类）",
        "每场定义核 + 资源导入标注",
        "带客、建群、回访摘要；园区做销售促成",
    ], BLUE)
    card(s, Inches(7.2), Inches(1.35), Inches(5.55), Inches(5.4), "请园区协同", [
        "确认档期、场地物业与样板间",
        "共享在谈名单，群内发布条件促成成交",
        "云创基地/联合会/服中心协同到位",
        "政策礼包与券务按活动节奏提前到位",
        "月度联席会核对意向与下月排期",
    ], GREEN)
    pg(s)

    # 22 thanks
    s = slide(prs)
    text(s, ML, Inches(2.3), Inches(12), Inches(0.5), "谢谢审阅，期待与您共事", size=32, color=INK, bold=True, align=PP_ALIGN.CENTER)
    text(s, ML, Inches(3.1), Inches(12), Inches(0.4), "创智汇年度活动运营方案 · 双会资源并入融合版", size=16, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    text(s, ML, Inches(3.65), Inches(12), Inches(0.35), "2026年8月启动 · 同浦汇带客 · 园区销售", size=14, color=MUT, align=PP_ALIGN.CENTER)
    text(s, ML, Inches(4.3), Inches(12), Inches(1.0),
         "主体运营：上海市云计算创新基地（国家级孵化器）\n学术支持：复旦大学住房政策研究中心\n载体支持：杨浦区科技企业联合会 · 科技企业服务中心",
         size=13, color=SOFT, align=PP_ALIGN.CENTER)
    pg(s)

    # 23-25 appendix definition maps
    for subset, title in ((EVENTS[:10], "附录 · 定义核与资源导入（1/3）"),
                          (EVENTS[10:20], "附录 · 定义核与资源导入（2/3）"),
                          (EVENTS[20:], "附录 · 定义核与资源导入（3/3）")):
        s = slide(prs); header(s, "附录", title, "每场一句话定义 + 导入什么资源（招商向，不堆全量）")
        rows = [["编号", "定义核", "资源源", "资源导入"]]
        for e in subset:
            rows.append([e["code"], e["define"][:22], e["src"][:14], e["res"][:28]])
        table(s, Inches(1.25), rows, [0.08, 0.32, 0.2, 0.4], sizes=[10, 11, 10, 10], rh=0.42)
        pg(s)

    # 26 parties one-pager (merged leftover)
    s = slide(prs); header(s, "附录 · 主体与口径一页纸", "合并后统一口径（便于审阅）")
    table(s, Inches(1.35), [
        ["模块", "口径"],
        ["主体运营范围", "上海市云计算创新基地（国家级孵化器）"],
        ["学术支持", "复旦大学住房政策研究中心"],
        ["活动及运营载体支持", "杨浦区科技企业联合会、科技企业服务中心"],
        ["执行带客", "同浦汇"],
        ["销售成交", "园区负责销售部分与合同交割"],
        ["人数/负责人", "每场≤30人；负责人约50%；全年约600人次+"],
        ["双会原则", "WAIC→3F；ChinaJoy→5F；具身/AIGC叠加；只取招商切片"],
        ["外企/出海/领事", "不承诺必带外企；出海另议；领事到访与挂牌均另计价且分性质"],
        ["租金与付款", "不对赌；免租期1–3个月面议；付款50%+4×10%+10%"],
    ], [0.28, 0.72], sizes=[12, 12], rh=0.42)
    pg(s)

    # 27-28 execution elements
    for subset, title in ((EVENTS[:15], "附录 · 逐场执行要素（上）"), (EVENTS[15:], "附录 · 逐场执行要素（下）")):
        s = slide(prs); header(s, "附录", title, "行程 / 嘉宾（与正文总表合并互补，不重复堆砌）")
        rows = [["编号", "活动", "行程", "嘉宾"]]
        for e in subset:
            rows.append([e["code"], e["name"][:14], e["agenda"][:26], e["guests"][:16]])
        table(s, Inches(1.25), rows, [0.08, 0.26, 0.38, 0.28], sizes=[9, 10, 10, 10], rh=0.32)
        pg(s)

    # 29 end note / merge说明
    s = slide(prs); header(s, "附录 · 合并说明", "本册如何把两套材料合成一套")
    card(s, ML, Inches(1.35), Inches(12.23), Inches(5.4), "合并同类项规则", [
        "结构：以《年度活动运营方案》章节为主骨架（背景→日历→六大主题→执行→转化收费案例→结语）",
        "场次：以修订后30场为准（≤30人/负责人约50%/8月项目推广日/ChinaJoy主题替换潮玩产业集群场）",
        "资源：把WAIC/ChinaJoy总表压缩为「可导入切片」写入主题页与附录定义核，不粘贴全量展商",
        "商务：付款唯一方案、不对赌、出海/领事另计价、带客/销售分工——统一进正文，去掉冲突旧表述",
        "案例：保留往期实拍案例页，映射到对应场次类型",
        "删除：关联度核验表交付、大场千人级承诺、租赁佣金对赌等与现行口径冲突项",
    ], BLUE)
    pg(s)

    assert n == TOTAL, f"page count {n} != {TOTAL}"

    os.makedirs(ART, exist_ok=True)
    outs = [
        "创智汇年度活动运营方案-双会并入融合版.pptx",
        "chuangzhihui-annual-ops-plan-dual-merge.pptx",
        "创智汇年度活动运营方案-融合版.pptx",
        "chuangzhihui-annual-ops-plan-fused.pptx",
    ]
    tmp = os.path.join(HERE, "_tmp_dual_merge.pptx")
    prs.save(tmp)
    for fn in outs:
        p = os.path.join(HERE, fn)
        shutil.copy2(tmp, p)
        shutil.copy2(p, os.path.join(ART, fn))
    os.remove(tmp)
    print("PPT slides:", len(prs.slides))

    # Excel
    wb = Workbook()
    thin = Border(left=Side(style="thin", color="D0D0D0"), right=Side(style="thin", color="D0D0D0"),
                  top=Side(style="thin", color="D0D0D0"), bottom=Side(style="thin", color="D0D0D0"))
    head = PatternFill("solid", fgColor="1E4D8C")
    alt = PatternFill("solid", fgColor="EEF2F8")

    def hdr(ws, headers):
        for j, h in enumerate(headers, 1):
            c = ws.cell(1, j, h); c.font = Font(name="微软雅黑", bold=True, color="FFFFFF", size=10)
            c.fill = head; c.alignment = Alignment(wrap_text=True, vertical="center"); c.border = thin

    def write(ws, data, widths, h=44):
        for i, row in enumerate(data, 2):
            for j, v in enumerate(row, 1):
                c = ws.cell(i, j, v); c.font = Font(name="微软雅黑", size=9)
                c.alignment = Alignment(wrap_text=True, vertical="center"); c.border = thin
                if i % 2 == 0: c.fill = alt
            ws.row_dimensions[i].height = h
        ws.row_dimensions[1].height = 28
        for i, w in enumerate(widths, 1):
            ws.column_dimensions[get_column_letter(i)].width = w

    ws = wb.active; ws.title = "01-活动概况"
    hdr(ws, ["编号", "活动名称", "时间", "形式", "人数", "负责人", "场地", "层次", "招商落点", "行程", "嘉宾"])
    write(ws, [[e["code"], e["name"], e["month"], e["fmt"], e["people"], e["dm"], e["floor"], e["layer"],
                e["lease"], e["agenda"], e["guests"]] for e in EVENTS],
          [8, 32, 16, 12, 10, 8, 12, 8, 22, 36, 28], 48)

    ws2 = wb.create_sheet("02-定义核资源导入")
    hdr(ws2, ["编号", "活动名称", "层次", "定义核", "资源源", "资源导入"])
    write(ws2, [[e["code"], e["name"], e["layer"], e["define"], e["src"], e["res"]] for e in EVENTS],
          [8, 32, 10, 28, 22, 40], 48)

    ws3 = wb.create_sheet("03-六大主题")
    hdr(ws3, ["主题", "名称", "场次编号", "说明"])
    write(ws3, [[t["key"], t["title"], "、".join(t["codes"]), "；".join(t["blurb"])] for t in THEMES],
          [8, 36, 28, 60], 56)

    ws4 = wb.create_sheet("04-主体口径收费")
    hdr(ws4, ["模块", "口径"])
    write(ws4, [
        ["主体运营范围", "上海市云计算创新基地（国家级孵化器）"],
        ["学术支持", "复旦大学住房政策研究中心"],
        ["载体支持", "杨浦区科技企业联合会、科技企业服务中心"],
        ["执行带客", "同浦汇"],
        ["销售", "园区负责销售部分"],
        ["人数", "每场≤30人；全年约600人次+"],
        ["负责人", "约50%"],
        ["双会", "WAIC→3F；ChinaJoy→5F；具身/AIGC叠加"],
        ["出海/领事", "另议另计价；到访≠挂牌"],
        ["租金", "3.3元/㎡/天；不做对赌与必要性要求；免租期1–3个月面议"],
        ["付款", "签约后7日50%；四季各10%；年终10%"],
        ["合并原则", "结构以年度方案PPT为主；场次与资源以双会修订版并入；删除冲突旧口径"],
    ], [18, 72], 34)

    xouts = [
        "创智汇年度活动运营方案-双会并入融合版.xlsx",
        "chuangzhihui-annual-ops-plan-dual-merge.xlsx",
        "创智汇年度活动运营方案-融合版.xlsx",
        "chuangzhihui-annual-ops-plan-fused.xlsx",
    ]
    for fn in xouts:
        p = os.path.join(HERE, fn); wb.save(p); shutil.copy2(p, os.path.join(ART, fn))
    print("Excel OK")


if __name__ == "__main__":
    build()
