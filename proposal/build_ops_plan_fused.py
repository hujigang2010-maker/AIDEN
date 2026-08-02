# -*- coding: utf-8 -*-
"""创智汇年度活动运营方案 · 融合版
以《创智汇年度活动运营方案（2026.08—2027.07）》PDF 为最终版式母版，
把既有 30 场交付（概况/价值/交付/分工/收费/案例照片/关联度口径）完整融合。
输出：PPT + PDF + Excel
用法: python3 build_ops_plan_fused.py
"""
import os, shutil
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
import fitz

HERE = os.path.dirname(os.path.abspath(__file__))
ART = "/opt/cursor/artifacts"
PDF_SRC = os.path.join(HERE, "创智汇年度活动运营方案2026.08-2027.07.pdf")
PHOTO = os.path.join(HERE, "cases", "photos")
ASSETS = os.path.join(HERE, "cases", "ops_assets")
RENDER = "/tmp/chuangzhihui_ops/hires"
FONT = "Microsoft YaHei"

# light theme matching PDF
INK = RGBColor(0x1A, 0x1A, 0x2E)
MUT = RGBColor(0x4A, 0x4A, 0x5A)
SOFT = RGBColor(0x8A, 0x8A, 0x9A)
BLUE = RGBColor(0x1E, 0x4D, 0x8C)
BLUE2 = RGBColor(0x2F, 0x6F, 0xB8)
GOLD = RGBColor(0xC4, 0x9A, 0x3C)
LINE = RGBColor(0xD8, 0xDE, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF7, 0xF8, 0xFA)
CARD = RGBColor(0xEE, 0xF2, 0xF8)
ACC_BG = RGBColor(0xE8, 0xF0, 0xFA)

SW, SH = Inches(13.333), Inches(7.5)  # matches 960x540
ML, CW = Inches(0.55), Inches(12.23)

# ---------- 30场（以PDF日历为准 2026.08—2027.07）+ 融合交付字段 ----------
EVENTS = [
    # 2026.08-12
    dict(code="A1", cat="A", name="智能伙伴·创智汇年度启动沙龙", month="2026年8月上旬", format="主题沙龙", people="60人", dm="80%",
         lease="集中邀约潜在客户，为全年定调", floor="5F沙龙",
         agenda="18:30签到参观→19:00年度主题解读→19:30园区定位发布→20:10圆桌→20:40看场意向",
         guests="联合会领导、AI/OPC企业负责人、服中心", value="开年定调；8间待租首波曝光",
         deliver="方案+议程+关联度核验表+签到+意向分级+7日回访+复盘"),
    dict(code="E1", cat="E", name="WAIC成果承接·创智汇开放日", month="2026年8月中旬", format="开放日", people="300人", dm="30%",
         lease="大会资源导入园区看场", floor="3F+5F整区",
         agenda="10:00致辞→10:30主题速递→11:30开放参观→14:00双展台→15:30政策礼包→16:30意向洽谈",
         guests="WAIC回流嘉宾、媒体、AI/内容企业、高校", value="大会流量导入；集中看场",
         deliver="同上+全园动线图+媒体通稿"),
    dict(code="E2", cat="E", name="AI Builders Night·建造者之夜", month="2026年8月中旬", format="夜场社交", people="120人", dm="50%",
         lease="建造者当场看3F单元", floor="5F主厅",
         agenda="签到名牌→致辞→闪电演讲×4→开放麦→AMA→投资人蹲项目+看场",
         guests="创业者、开发者、投资人、开源Maintainer", value="Builders当场看单元",
         deliver="同上+项目对接表"),
    dict(code="E3", cat="E", name="通往AGI之夜", month="2026年8月下旬", format="学术夜场", people="80人", dm="60%",
         lease="招引高端涉外与研究型企业", floor="5F圆桌区",
         agenda="复旦致辞→无PPT圆桌→三桌围炉→辩题投票→深度社交+精品看场",
         guests="复旦学者、大模型负责人、Agent创始人、治理学者", value="高端研究型客户",
         deliver="同上"),
    dict(code="C5", cat="C", name="东莞潮玩品牌入沪推介", month="2026年8月下旬", format="推介会", people="50人", dm="80%",
         lease="品牌展位与快闪合作", floor="5F潮玩区",
         agenda="上海渠道→5F落位→联名活动→看场→定金意向",
         guests="潮玩品牌、渠道商负责人", value="品牌展位/快闪去化",
         deliver="同上+报价单"),
    dict(code="A2", cat="A", name="Agent智能体搭建一日营", month="2026年9月上旬", format="实训营", people="40人", dm="60%",
         lease="推动Agent初创团队落位", floor="3F OPC",
         agenda="09:30签到→10:00工作流拆解→13:30低代码实操→16:00Demo→16:40看场逼单",
         guests="Agent工程师、开源Maintainer、初创CEO", value="Agent团队入孵",
         deliver="同上+作品集"),
    dict(code="C6", cat="C", name="IP授权×AI衍生交易撮合会", month="2026年9月中旬", format="撮合会", people="60人", dm="80%",
         lease="IP与衍生品团队入驻", floor="5F展示中心",
         agenda="IP路演→AI衍生专题桌→一对一撮合→成交看板→落户激励",
         guests="IP方、被授权商、律师、渠道", value="IP/衍生品入驻；贸易成交",
         deliver="同上+成交登记"),
    dict(code="F1", cat="F", name="AI出海推介与国际买家团", month="2026年9月下旬", format="买家团", people="60人", dm="85%",
         lease="出海需求企业入驻", floor="5F+培训",
         agenda="出海路径→买家需求→一对一洽谈→订单意向→办公/展位捆绑",
         guests="海外买家、跨境平台、出海企业负责人", value="出海入驻；撮合收入",
         deliver="同上"),
    dict(code="A3", cat="A", name="多模态智能体工作坊", month="2026年10月上旬", format="工作坊", people="40人", dm="60%",
         lease="吸引多模态应用团队", floor="3F培训+直播间",
         agenda="案例分享→视觉语言实操→场景共创→作品评审→政策礼包+看场",
         guests="多模态研究者、应用厂商产品负责人", value="多模态团队入驻",
         deliver="同上"),
    dict(code="E4", cat="E", name="创智汇秋季潮玩×AI嘉年华", month="2026年10月中旬", format="嘉年华（2天）", people="500人", dm="25%",
         lease="集中签约展位与办公", floor="5F+公区",
         agenda="开幕秀→市集展洽→AI赛事→夜间场→签约仪式",
         guests="潮玩品牌、达人、赞助商、AI创作者", value="集中签约",
         deliver="同上+赞助结算"),
    dict(code="B2", cat="B", name="AI治理与可信智能体沙龙", month="2026年10月下旬", format="沙龙", people="40人", dm="85%",
         lease="建立合规型企业信任", floor="3F",
         agenda="治理框架→合规清单→案例→问答→入驻服务对接",
         guests="治理学者、合规顾问、企业法务/CEO", value="合规型企业入驻",
         deliver="同上"),
    dict(code="A4", cat="A", name="火山引擎×算力Infra实务营", month="2026年11月上旬", format="厂商联训", people="40人", dm="70%",
         lease="定向招引高算力企业", floor="3F+云创基地",
         agenda="算力政策包→代金券实操→集群案例→一对一诊断→入驻捆绑云资源",
         guests="火山引擎讲师、云创基地、高算力企业CTO", value="高算力企业定向招商",
         deliver="同上"),
    dict(code="B4", cat="B", name="创新券·算力券·模型券实务沙龙", month="2026年11月中旬", format="实务沙龙", people="50人", dm="85%",
         lease="促进用券企业向园区聚集", floor="3F",
         agenda="三券规则→核销实操→案例→开户引导→入驻转化",
         guests="创新券平台、云厂商、用券企业负责人", value="用券企业聚集",
         deliver="同上"),
    dict(code="F2", cat="F", name="国别会客厅·AI合作领事专题", month="2026年11月下旬", format="领事沙龙", people="40人", dm="90%",
         lease="链接涉外与出海企业", floor="会客厅/5F",
         agenda="外事接待→国别AI机会→B2B闪见→国际服务→高端看场",
         guests="领事官员、涉外AI企业负责人", value="涉外/出海入驻背书",
         deliver="同上"),
    dict(code="A5", cat="A", name="具身智能空间交互体验日", month="2026年12月上旬", format="体验日", people="60人", dm="50%",
         lease="组织机器人团队看场", floor="5F展区+3F",
         agenda="具身演示→空间交互讲解→场景清单→5F踩点→落位洽谈",
         guests="机器人/具身团队、高校实验室", value="具身团队看场",
         deliver="同上"),
    dict(code="F4", cat="F", name="创智汇AI年度Demo Day", month="2026年12月中旬", format="路演日", people="200人", dm="65%",
         lease="集中签约与媒体背书", floor="5F主展",
         agenda="开幕→十强路演→颁奖→明年预热→红酒洽谈+集中签约",
         guests="投资人、链主、媒体、区级嘉宾、十强项目", value="集中签约；媒体背书",
         deliver="同上+签约台账"),
    # 2027.01-07
    dict(code="D1", cat="D", name="AIGC微短剧制片特训", month="2027年1月上旬", format="特训营", people="40人", dm="60%",
         lease="厂牌与工作室入驻", floor="3F直播间+5F",
         agenda="YOUNG立方政策→制片流程→脚本工作坊→路演→5F看场",
         guests="微短剧导演、厂牌负责人、平台方", value="内容厂牌入驻",
         deliver="同上"),
    dict(code="A6", cat="A", name="AI营销Agent实战营", month="2027年1月中旬", format="实战营", people="40人", dm="70%",
         lease="招引营销科技类公司", floor="3F培训",
         agenda="投放策略→Agent自动化→素材AIGC→复盘→服务包转化",
         guests="投放操盘手、Agent产品经理", value="营销科技入驻",
         deliver="同上"),
    dict(code="B1", cat="B", name="YOUNG立方×智能伙伴政策沙龙", month="2027年2月中旬", format="政策沙龙", people="50人", dm="85%",
         lease="导入内容与AI企业", floor="5F沙龙",
         agenda="政策要点→园区礼包→适用画像→诊断预约→看场",
         guests="区政策宣讲、服中心、内容/AI企业负责人", value="政策前置获客",
         deliver="同上"),
    dict(code="B3", cat="B", name="高企认定冲刺（AI企业专场）", month="2027年2月下旬", format="辅导会", people="40人", dm="90%",
         lease="外区待认定企业带政策入驻", floor="3F",
         agenda="条件拆解→材料要点→初筛→辅导套餐→入驻激励",
         guests="高企辅导、财税顾问、待认定AI企业负责人", value="带政策入驻",
         deliver="同上"),
    dict(code="A7", cat="A", name="OPC超级个体黑客松（春）", month="2027年3月上旬", format="黑客松（1.5天）", people="80人", dm="45%",
         lease="获奖团队优先洽谈单元", floor="3F整层",
         agenda="开题组队→封闭开发→路演决赛→入驻礼包→一对一看场",
         guests="评委×3、投资人×2、云厂商、Builders", value="获奖团队谈单元",
         deliver="同上"),
    dict(code="C2", cat="C", name="高校成果转化·AI for Science日", month="2027年3月下旬", format="对接日", people="60人", dm="80%",
         lease="成果公司与实验室落户", floor="3F+沙龙",
         agenda="成果路演→科学智能体场景→园区承接→看孵化单元→转化落位",
         guests="复旦/同济技术转移、教授团队", value="成果公司落户",
         deliver="同上"),
    dict(code="F3", cat="F", name="通往AGI季度圆桌", month="2027年3月下旬", format="闭门圆桌", people="30人", dm="90%",
         lease="研究型与模型团队", floor="5F沙龙",
         agenda="议题导入→圆桌→围炉→问答→看场",
         guests="学者、模型企业、投资人", value="研究型/模型团队",
         deliver="同上"),
    dict(code="D2", cat="D", name="AI创作者内容首发①", month="2027年4月中旬", format="发布会", people="100人", dm="40%",
         lease="带动关联品牌问询", floor="5F主展",
         agenda="揭幕→发布→签售快闪→媒体专访→招商通道",
         guests="创作者、媒体、渠道、品牌负责人", value="品牌问询",
         deliver="同上"),
    dict(code="C3", cat="C", name="汕头玩具×AIGC供应链对接会", month="2027年4月下旬", format="对接会", people="50人", dm="80%",
         lease="玩具企业展位落位", floor="5F玩具区",
         agenda="集群介绍→AIGC设计提效→供需对接→展位参观→报价",
         guests="汕头商协会、品牌采购、玩具企业负责人", value="玩具展位去化",
         deliver="同上"),
    dict(code="C1", cat="C", name="专精特新·AI应用培育路演", month="2027年5月中旬", format="路演", people="80人", dm="75%",
         lease="补位成长型AI企业", floor="5F",
         agenda="梯度政策→企业路演→点评→金融对接→看场",
         guests="评审顾问、银行、基金、成长型AI企业", value="成长型AI补位",
         deliver="同上"),
    dict(code="D3", cat="D", name="潮玩×AIGC主题市集", month="2027年5月下旬", format="市集（2天）", people="400人", dm="25%",
         lease="摊主升级固定展位", floor="5F+公区",
         agenda="布摊→开市→AIGC打卡→夜间场→优质摊主转正谈",
         guests="摊主×40、达人、潮玩品牌", value="摊主转正/办公",
         deliver="同上"),
    dict(code="C4", cat="C", name="扬州毛绒×数字人联名对接", month="2027年6月中旬", format="对接会", people="40人", dm="85%",
         lease="700㎡量级客户", floor="5F毛绒区",
         agenda="联名模式→租金扣点测算→看场→MOU→跟进",
         guests="扬州集群、数字人厂商、品牌方", value="700㎡量级客户",
         deliver="同上"),
    dict(code="D4", cat="D", name="沉浸式AI+IP联展", month="2027年6月下旬", format="联展（7–10天）", people="2000人次", dm="15%",
         lease="闭幕酒会集中转化", floor="5F展示中心",
         agenda="布展→开幕→预约观展→教育场→闭幕招商酒会",
         guests="联合IP×5、文旅渠道、品牌", value="闭幕酒会集中转化",
         deliver="同上+展期日报"),
    dict(code="D5", cat="D", name="AI创作者内容首发②（WAIC2027预热场）", month="2027年7月中旬", format="发布会", people="100人", dm="40%",
         lease="旺季补位，衔接新一届大会", floor="5F",
         agenda="发布→渠道对接→看场→年框意向→WAIC2027预热",
         guests="IP方、渠道、媒体", value="衔接WAIC2027",
         deliver="同上"),
]

assert len(EVENTS) == 30

CASES = [
    {
        "title": "往期案例（一）训练营 / 沙龙类",
        "name": "杨「数」浦数字沙龙第七期：AI如何重塑企业DNA",
        "meta": "2025-06-30 · 美团上海综合指挥中心 · 政策/治理沙龙类",
        "orgs": "杨浦区委网信办主办；赛博网络安全产业创新研究院、上海市人工智能与社会发展研究会承办",
        "companies": "美团；大众点评大模型应用团队；赛博院研究员卢虹羽等",
        "agenda": [
            ("参观", "美团综合指挥中心数字化场景参观"),
            ("主题", "人机共生时代：AI如何重塑企业DNA"),
            ("案例", "大众点评AI技术应用创新与落地"),
            ("政策", "《人工智能生成合成内容标识办法》等解读"),
            ("互动", "答疑与企业一对一交流"),
        ],
        "photos": ["case_c987bd91.jpg", "case_59157ea5.jpg", "case_6e51af9e.jpg"],
        "result": "面向数字经济企业管理层/法务/技术；公开报道见杨浦区政府官网。对应创智汇：B2治理沙龙、B1政策沙龙、A类训练营。",
    },
    {
        "title": "往期案例（二）开放日 / 大场 / 路演类",
        "name": "「融见科创·智启未来」人工智能专场路演暨投融资对接会",
        "meta": "2025-10 · 杨浦 · 邮储银行联合主办 · 路演大场类",
        "orgs": "杨浦科创促进会、邮储银行上海分行；科辰创投、智能产业创新研究院、益华资本等",
        "companies": "复楚智能、卡房信息、一造科技、中科趋势、万笔千墨等路演企业",
        "agenda": [
            ("开场", "促进会会长夏立城致辞与生态介绍"),
            ("主办", "邮储银行科技金融致辞"),
            ("分享", "中邮证券AI行业趋势主题分享"),
            ("路演", "“6+6+6+1”模式：6项目+6投资人+6金融机构+1政府代表"),
            ("收口", "区投促办点评与招商邀请；评委总结跟进"),
        ],
        "photos": ["case_db1d5cac.jpg", "case_65931516.jpg", "case_cb0a3812.jpg"],
        "result": "可见企业路演+金融机构对接。对应创智汇：E1开放日、E4嘉年华、F4 Demo Day、C1专精特新路演。另附杨数浦出海合规授证现场作为出海/大场仪式参考。",
    },
]


def ensure_renders():
    os.makedirs(RENDER, exist_ok=True)
    need = [i for i in range(1, 28) if not os.path.exists(f"{RENDER}/p{i:02d}.png")]
    if not need:
        return
    doc = fitz.open(PDF_SRC)
    for i in need:
        pix = doc[i - 1].get_pixmap(matrix=fitz.Matrix(2.2, 2.2))
        pix.save(f"{RENDER}/p{i:02d}.png")
    doc.close()


def _set_run(p, text, size, color, bold=False):
    r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = FONT
    return r


def textbox(s, x, y, w, h, text, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]; p.alignment = align
    _set_run(p, text, size, color, bold)
    return tb


def rect(s, x, y, w, h, fill=None, line=None, lw=1.0):
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    b.shadow.inherit = False
    if fill is None: b.fill.background()
    else: b.fill.solid(); b.fill.fore_color.rgb = fill
    if line is None: b.line.fill.background()
    else: b.line.color.rgb = line; b.line.width = Pt(lw)
    return b


def add_pic(s, path, x, y, w, h):
    if os.path.exists(path):
        s.shapes.add_picture(path, x, y, width=w, height=h)
    else:
        rect(s, x, y, w, h, fill=CARD, line=LINE)
        textbox(s, x, y, w, h, "照片待补", size=12, color=SOFT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def light_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, SW, SH, fill=WHITE)
    return s


def footer_bar(s, page, total):
    textbox(s, ML, Inches(7.12), Inches(8), Inches(0.28),
            "上海创智汇 × 同浦汇 · 2026.08—2027.07 · 融合版", size=10, color=SOFT)
    textbox(s, Inches(11.2), Inches(7.12), Inches(1.6), Inches(0.28),
            f"{page} / {total}", size=10, color=SOFT, align=PP_ALIGN.RIGHT)


def header_bar(s, sec, title, sub=None):
    rect(s, 0, 0, SW, Inches(0.08), fill=BLUE)
    textbox(s, ML, Inches(0.28), Inches(12), Inches(0.28), sec, size=12, color=BLUE, bold=True)
    textbox(s, ML, Inches(0.55), Inches(12), Inches(0.45), title, size=22, color=INK, bold=True)
    if sub:
        textbox(s, ML, Inches(1.0), Inches(12), Inches(0.3), sub, size=12, color=MUT)


def table_like(s, y0, rows, col_w, sizes=None, rh=0.34, head=True):
    yy = y0
    if sizes is None: sizes = [10] * len(rows[0])
    for ri, row in enumerate(rows):
        xx = ML
        h = Inches(0.4 if ri == 0 and head else rh)
        for ci, val in enumerate(row):
            cw = Emu(int(CW * col_w[ci]))
            bg = BLUE if ri == 0 and head else (ACC_BG if ri % 2 == 0 else WHITE)
            fg = WHITE if ri == 0 and head else INK
            rect(s, xx, yy, cw, h, fill=bg, line=LINE, lw=0.5)
            textbox(s, Emu(xx + Inches(0.05)), Emu(yy + Inches(0.04)), Emu(cw - Inches(0.08)), Emu(h - Inches(0.06)),
                    str(val), size=sizes[ci], color=fg, bold=(ri == 0 and head), anchor=MSO_ANCHOR.MIDDLE)
            xx = Emu(xx + cw)
        yy = Emu(yy + h)
    return yy


def build_ppt():
    ensure_renders()
    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH

    # pages 1-23, 26-27 from PDF renders (exact format)
    # pages 24-25 rebuilt with filled cases
    # then appendix pages with fused detail
    base_pages = list(range(1, 24)) + [26, 27]  # skip 24,25 for rebuild
    # We'll insert: 1..23 (pdf), 24-25 (cases), 26-27 (pdf), then appendix

    for i in range(1, 24):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_pic(s, f"{RENDER}/p{i:02d}.png", 0, 0, SW, SH)

    # ---- rebuilt case pages 24-25 ----
    # Case 1
    s = light_slide(prs)
    header_bar(s, "0 5 · 往期案例（一）训练营 / 沙龙类", CASES[0]["name"], CASES[0]["meta"])
    c = CASES[0]
    # photos
    for i, fn in enumerate(c["photos"][:3]):
        add_pic(s, os.path.join(PHOTO, fn), Emu(ML + i * Inches(4.1)), Inches(1.4), Inches(3.95), Inches(2.35))
    rect(s, ML, Inches(3.9), Inches(7.6), Inches(2.9), fill=BG, line=LINE)
    textbox(s, Emu(ML + Inches(0.2)), Inches(4.0), Inches(7.2), Inches(0.3), "当天议程", size=13, color=BLUE, bold=True)
    yy = Inches(4.35)
    for t, d in c["agenda"]:
        textbox(s, Emu(ML + Inches(0.2)), yy, Inches(1.0), Inches(0.28), t, size=11, color=GOLD, bold=True)
        textbox(s, Emu(ML + Inches(1.2)), yy, Inches(6.0), Inches(0.28), d, size=11, color=INK)
        yy = Emu(yy + Inches(0.28))
    rect(s, Inches(8.4), Inches(3.9), Inches(4.4), Inches(2.9), fill=ACC_BG, line=LINE)
    textbox(s, Inches(8.55), Inches(4.0), Inches(4.1), Inches(0.3), "到场机构 / 企业", size=13, color=BLUE, bold=True)
    textbox(s, Inches(8.55), Inches(4.4), Inches(4.1), Inches(1.0), f"{c['orgs']}\n\n企业侧：{c['companies']}", size=11, color=INK)
    textbox(s, Inches(8.55), Inches(5.5), Inches(4.1), Inches(1.1), c["result"], size=11, color=MUT)
    footer_bar(s, 24, 32)

    # Case 2
    s = light_slide(prs)
    header_bar(s, "0 5 · 往期案例（二）开放日 / 嘉年华 / 路演大场类", CASES[1]["name"], CASES[1]["meta"])
    c = CASES[1]
    add_pic(s, os.path.join(PHOTO, c["photos"][0]), ML, Inches(1.4), Inches(6.0), Inches(2.5))
    add_pic(s, os.path.join(PHOTO, c["photos"][1]), Inches(6.75), Inches(1.4), Inches(2.95), Inches(2.5))
    add_pic(s, os.path.join(PHOTO, c["photos"][2]), Inches(9.9), Inches(1.4), Inches(2.85), Inches(2.5))
    rect(s, ML, Inches(4.1), Inches(7.8), Inches(2.7), fill=BG, line=LINE)
    textbox(s, Emu(ML + Inches(0.2)), Inches(4.2), Inches(7.4), Inches(0.3), "当天议程与成效", size=13, color=BLUE, bold=True)
    yy = Inches(4.55)
    for t, d in c["agenda"]:
        textbox(s, Emu(ML + Inches(0.2)), yy, Inches(1.0), Inches(0.28), t, size=11, color=GOLD, bold=True)
        textbox(s, Emu(ML + Inches(1.2)), yy, Inches(6.2), Inches(0.28), d, size=11, color=INK)
        yy = Emu(yy + Inches(0.28))
    rect(s, Inches(8.55), Inches(4.1), Inches(4.25), Inches(2.7), fill=ACC_BG, line=LINE)
    textbox(s, Inches(8.7), Inches(4.2), Inches(4.0), Inches(0.3), "可见企业 / 映射", size=13, color=BLUE, bold=True)
    textbox(s, Inches(8.7), Inches(4.55), Inches(4.0), Inches(2.1),
            f"主办协同：{c['orgs']}\n\n路演企业：{c['companies']}\n\n{c['result']}", size=11, color=INK)
    footer_bar(s, 25, 32)

    # pages 26-27 from PDF
    for i in (26, 27):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_pic(s, f"{RENDER}/p{i:02d}.png", 0, 0, SW, SH)

    # ========== 融合附录 ==========
    # A1 关联度口径（对齐PDF：100%定向邀约）
    s = light_slide(prs)
    header_bar(s, "附录 A · 客群与孵化器关联度", "定向邀约 100% 主题关联——可核验口径",
               "对齐本方案总表脚注：全部场次不设随机散客渠道")
    rows = [
        ["口径", "标准", "核验方式", "交付物"],
        ["邀约名单", "100% 按园区产业主题定向邀约", "报名表行业标签审核", "关联度核验表"],
        ["散客渠道", "不设随机散客 / 公开路人渠道", "报名通道仅定向名单", "报名页权限说明"],
        ["负责人占比", "按场次目标值（总表列示）", "职位字段统计", "单场复盘"],
        ["生态角色", "政府/高校/投资人/媒体为背书席位", "单独标注，不计入散客", "嘉宾构成表"],
        ["月度公示", "联席会核对达成率", "台账对账", "转化月报"],
    ]
    table_like(s, Inches(1.45), rows, [0.16, 0.32, 0.28, 0.24], sizes=[11, 12, 12, 12], rh=0.55)
    textbox(s, ML, Inches(5.2), CW, Inches(1.5),
            "说明：本融合版以运营方案 PDF 总表脚注为准——客群与孵化器方向关联度 100%（定向邀约）。"
            "负责人占比为运营目标值，签约后按客群名单逐场确认。每场活动结束后 7 天内归档《关联度核验表》。",
            size=12, color=MUT)
    footer_bar(s, 28, 32)

    # A2 价值·交付·分工·收费一页总览（融合原表格应答）
    s = light_slide(prs)
    header_bar(s, "附录 B · 融合总览", "价值 · 交付 · 分工 · 收费（一张表说清）")
    rows = [
        ["模块", "内容"],
        ["活动价值", "议题变日常；3F/5F主题坐实；载体考核与补贴台账；看场≥20分钟/场；7日三轮跟进；导流8间待租"],
        ["具体交付", "方案/议程/嘉宾确认函/报名页/签到/意向表/看场动线/回访/报价/复盘/照片通稿；月报；季报；年报"],
        ["我方主责", "策划、定向邀约、物料、执行、看场带队、回访报价、台账复盘、传播"],
        ["园区支持", "方案审定、档期、在谈名单共享、场地物业导视安保、样板间、价格审批、入驻交割"],
        ["服中心/高校", "券务受理与申报辅导、成果团队与学者推荐、联合会渠道联动"],
        ["收费模式", "活动策划执行服务费另行商定（可年度打包/按场）；佣金=成交首年有效净租金×2个月；租金3.3元/㎡/天；物业13.8元/㎡/月"],
    ]
    table_like(s, Inches(1.4), rows, [0.16, 0.84], sizes=[12, 12], rh=0.62)
    footer_bar(s, 29, 32)

    # A3 上半年融合明细（行程/嘉宾）
    s = light_slide(prs)
    header_bar(s, "附录 C · 逐场执行要素（上）", "2026.08—12 · 行程 / 嘉宾 / 交付（融合原30场执行卡）")
    rows = [["编号", "活动", "行程要点", "嘉宾", "交付"]]
    for e in EVENTS[:16]:
        rows.append([e["code"], e["name"][:12], e["agenda"][:28], e["guests"][:16], "标准包"])
    table_like(s, Inches(1.35), rows, [0.07, 0.2, 0.35, 0.26, 0.12], sizes=[9, 9, 9, 9, 9], rh=0.3)
    footer_bar(s, 30, 32)

    # A4 下半年
    s = light_slide(prs)
    header_bar(s, "附录 C · 逐场执行要素（下）", "2027.01—07 · 行程 / 嘉宾 / 交付（融合原30场执行卡）")
    rows = [["编号", "活动", "行程要点", "嘉宾", "交付"]]
    for e in EVENTS[16:]:
        rows.append([e["code"], e["name"][:12], e["agenda"][:28], e["guests"][:16], "标准包"])
    table_like(s, Inches(1.35), rows, [0.07, 0.2, 0.35, 0.26, 0.12], sizes=[9, 9, 9, 9, 9], rh=0.32)
    footer_bar(s, 31, 32)

    # A5 WAIC映射 + 结语
    s = light_slide(prs)
    header_bar(s, "附录 D · WAIC 议程映射与融合说明", "大会一周 → 园区一年（融合 WAIC 对齐版要点）")
    rows = [
        ["WAIC/WAICA赛道", "创智汇承接", "代表场次"],
        ["智能伙伴主旨", "年度启动沙龙 + 开放日", "A1 / E1"],
        ["Agent智能体", "搭建营 / 营销Agent营", "A2 / A6"],
        ["多模态 / AI for Science", "工作坊 / 高校成果日", "A3 / C2"],
        ["算力Infra", "火山引擎实务营 / 三券沙龙", "A4 / B4"],
        ["具身智能", "空间交互体验日", "A5"],
        ["Builders / 青年菁英", "Builders Night / 黑客松", "E2 / A7"],
        ["AGI / 治理", "AGI之夜 / 治理沙龙 / 圆桌", "E3 / B2 / F3"],
        ["内容创作（错峰）", "特训 / 首发 / 市集 / 联展", "D1–D5"],
        ["出海 / 国际", "买家团 / 领事专题", "F1 / F2"],
    ]
    table_like(s, Inches(1.35), rows, [0.28, 0.42, 0.3], sizes=[11, 12, 11], rh=0.38)
    textbox(s, ML, Inches(6.0), CW, Inches(0.8),
            "融合说明：正文 27 页版式与日历完全以贵方上传的《创智汇年度活动运营方案（2026.08—2027.07）》为准；"
            "附录补入原30场交付中的执行卡要素、关联度核验口径、案例实拍与 WAIC 映射，形成可执行的一套材料。",
            size=11, color=MUT)
    footer_bar(s, 32, 32)

    # save
    os.makedirs(ART, exist_ok=True)
    outs = [
        "创智汇年度活动运营方案-融合版.pptx",
        "chuangzhihui-annual-ops-plan-fused.pptx",
        "创智汇30场活动具体方案.pptx",  # overwrite main events deliverable with fused master
        "chuangzhihui-30-events-plan.pptx",
    ]
    tmp = os.path.join(HERE, "_tmp_fused.pptx")
    prs.save(tmp)
    for fn in outs:
        p = os.path.join(HERE, fn)
        shutil.copy2(tmp, p)
        shutil.copy2(p, os.path.join(ART, fn))
    os.remove(tmp)
    print("PPT slides:", len(prs.slides))
    return len(prs.slides)


def build_pdf_from_ppt_pages():
    """Also produce a fused PDF: PDF pages 1-23 + new case renders + 26-27 + appendix note pages.
    Simpler: convert fused PPT pages by reusing PDF + inserting case pages via pymupdf.
    """
    # Build PDF by composing: original pages with 24-25 replaced, plus appendix pages rendered from PPT is hard.
    # Instead: create PDF from hi-res images of original + newly rendered case/appendix via fitz from a second pass.
    # Practical approach: export PDF that is original with replaced 24-25 and appended text pages for appendix.

    doc = fitz.open(PDF_SRC)
    # replace pages 24 and 25 by inserting images of our case slides - we need to render them.
    # Easier: create new PDF from scratch using original page pixmaps + new pages.

    out = fitz.open()
    # pages 1-23
    for i in range(23):
        out.insert_pdf(doc, from_page=i, to_page=i)

    # create case pages as PDF pages with images
    def add_image_page(img_path):
        # create page 960x540
        page = out.new_page(width=960, height=540)
        # we'll build case pages differently - use a white page with text for PDF appendix
        return page

    # For case pages, render using temporary images created below
    case_dir = "/tmp/fused_case_pages"
    os.makedirs(case_dir, exist_ok=True)
    make_case_page_images(case_dir)

    for fn in ["case1.png", "case2.png"]:
        path = os.path.join(case_dir, fn)
        page = out.new_page(width=960, height=540)
        page.insert_image(page.rect, filename=path)

    # pages 26-27 from original
    out.insert_pdf(doc, from_page=25, to_page=26)

    # appendix pages as images
    for fn in ["appA.png", "appB.png", "appC1.png", "appC2.png", "appD.png"]:
        path = os.path.join(case_dir, fn)
        if os.path.exists(path):
            page = out.new_page(width=960, height=540)
            page.insert_image(page.rect, filename=path)

    pdf_out = os.path.join(HERE, "创智汇年度活动运营方案-融合版.pdf")
    pdf_out_en = os.path.join(HERE, "chuangzhihui-annual-ops-plan-fused.pdf")
    out.save(pdf_out)
    shutil.copy2(pdf_out, pdf_out_en)
    shutil.copy2(pdf_out, os.path.join(ART, os.path.basename(pdf_out)))
    shutil.copy2(pdf_out_en, os.path.join(ART, os.path.basename(pdf_out_en)))
    # also overwrite the uploaded-name style deliverable
    shutil.copy2(pdf_out, os.path.join(HERE, "创智汇年度活动运营方案（2026.08—2027.07）-融合版.pdf"))
    shutil.copy2(pdf_out, os.path.join(ART, "创智汇年度活动运营方案（2026.08—2027.07）-融合版.pdf"))
    print("PDF pages:", out.page_count)
    out.close(); doc.close()


def make_case_page_images(outdir):
    """Render case+appendix pages to PNG via PIL for PDF embedding."""
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        import subprocess
        subprocess.check_call(["pip", "install", "pillow", "-q"])
        from PIL import Image, ImageDraw, ImageFont

    W, H = 1920, 1080  # 2x 960x540

    def font(size, bold=False):
        for path in [
            "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
            "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
            "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc",
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        ]:
            if os.path.exists(path):
                try:
                    return ImageFont.truetype(path, size)
                except Exception:
                    continue
        return ImageFont.load_default()

    def new_img():
        im = Image.new("RGB", (W, H), (255, 255, 255))
        d = ImageDraw.Draw(im)
        d.rectangle([0, 0, W, 8], fill=(30, 77, 140))
        return im, d

    def paste_photo(im, path, box):
        if not os.path.exists(path):
            return
        p = Image.open(path).convert("RGB")
        p = p.resize((box[2] - box[0], box[3] - box[1]), Image.Resampling.LANCZOS)
        im.paste(p, box[:2])

    # Case 1
    im, d = new_img()
    d.text((80, 40), "05 · 往期案例（一）训练营 / 沙龙类", fill=(30, 77, 140), font=font(28))
    d.text((80, 90), CASES[0]["name"], fill=(26, 26, 46), font=font(36))
    d.text((80, 150), CASES[0]["meta"], fill=(74, 74, 90), font=font(22))
    photos = CASES[0]["photos"]
    for i, fn in enumerate(photos[:3]):
        x0 = 80 + i * 600
        paste_photo(im, os.path.join(PHOTO, fn), (x0, 200, x0 + 560, 520))
    d.rectangle([80, 560, 1180, 1000], fill=(247, 248, 250), outline=(216, 222, 232))
    d.text((100, 580), "当天议程", fill=(30, 77, 140), font=font(26))
    y = 630
    for t, txt in CASES[0]["agenda"]:
        d.text((100, y), t, fill=(196, 154, 60), font=font(22))
        d.text((220, y), txt, fill=(26, 26, 46), font=font(22))
        y += 50
    d.rectangle([1220, 560, 1840, 1000], fill=(232, 240, 250), outline=(216, 222, 232))
    d.text((1240, 580), "到场机构 / 企业", fill=(30, 77, 140), font=font(26))
    d.text((1240, 640), CASES[0]["orgs"][:48], fill=(26, 26, 46), font=font(20))
    d.text((1240, 720), "企业侧：" + CASES[0]["companies"][:40], fill=(26, 26, 46), font=font(20))
    d.text((1240, 820), CASES[0]["result"][:70], fill=(74, 74, 90), font=font(18))
    d.text((80, 1020), "上海创智汇 × 同浦汇 · 2026.08—2027.07 · 融合版", fill=(138, 138, 154), font=font(18))
    d.text((1700, 1020), "24 / 32", fill=(138, 138, 154), font=font(18))
    im.save(os.path.join(outdir, "case1.png"))

    # Case 2
    im, d = new_img()
    d.text((80, 40), "05 · 往期案例（二）开放日 / 嘉年华 / 路演大场类", fill=(30, 77, 140), font=font(28))
    d.text((80, 90), CASES[1]["name"], fill=(26, 26, 46), font=font(32))
    d.text((80, 145), CASES[1]["meta"], fill=(74, 74, 90), font=font(22))
    paste_photo(im, os.path.join(PHOTO, CASES[1]["photos"][0]), (80, 200, 900, 560))
    paste_photo(im, os.path.join(PHOTO, CASES[1]["photos"][1]), (940, 200, 1380, 560))
    paste_photo(im, os.path.join(PHOTO, CASES[1]["photos"][2]), (1420, 200, 1840, 560))
    d.rectangle([80, 590, 1180, 1000], fill=(247, 248, 250), outline=(216, 222, 232))
    d.text((100, 610), "当天议程", fill=(30, 77, 140), font=font(26))
    y = 660
    for t, txt in CASES[1]["agenda"]:
        d.text((100, y), t, fill=(196, 154, 60), font=font(22))
        d.text((220, y), txt[:42], fill=(26, 26, 46), font=font(20))
        y += 48
    d.rectangle([1220, 590, 1840, 1000], fill=(232, 240, 250), outline=(216, 222, 232))
    d.text((1240, 610), "可见企业 / 映射", fill=(30, 77, 140), font=font(26))
    d.text((1240, 670), CASES[1]["companies"][:50], fill=(26, 26, 46), font=font(20))
    d.text((1240, 780), CASES[1]["result"][:80], fill=(74, 74, 90), font=font(18))
    d.text((80, 1020), "上海创智汇 × 同浦汇 · 2026.08—2027.07 · 融合版", fill=(138, 138, 154), font=font(18))
    d.text((1700, 1020), "25 / 32", fill=(138, 138, 154), font=font(18))
    im.save(os.path.join(outdir, "case2.png"))

    # Appendix pages (simplified text boards)
    apps = [
        ("appA.png", "附录 A · 客群与孵化器关联度",
         ["邀约名单：100% 按园区产业主题定向邀约（对齐总表脚注）",
          "不设随机散客 / 公开路人渠道",
          "负责人占比：按场次目标值，签约后逐场确认",
          "生态角色（政府/高校/投资人/媒体）单独标注为背书席位",
          "每场交付《关联度核验表》，月度联席会公示达成率"]),
        ("appB.png", "附录 B · 价值 · 交付 · 分工 · 收费",
         ["价值：议题日常化；3F/5F主题坐实；载体考核台账；8间待租导流",
          "交付：方案/议程/嘉宾函/报名/签到/意向/动线/回访/报价/复盘/通稿；月季年报",
          "我方：策划邀约物料执行看场回访台账传播",
          "园区：审定档期名单场地物业样板间价格审批入驻交割",
          "收费：活动服务费另议；佣金=2个月租金；办公3.3元/㎡/天；物业13.8元/㎡/月"]),
        ("appC1.png", "附录 C · 逐场执行要素（2026.08—12）",
         [f"{e['code']} {e['name']}｜{e['agenda'][:36]}" for e in EVENTS[:10]]),
        ("appC2.png", "附录 C · 逐场执行要素（续）",
         [f"{e['code']} {e['name']}｜{e['agenda'][:36]}" for e in EVENTS[10:20]]),
        ("appD.png", "附录 D · WAIC 映射与融合说明",
         ["智能伙伴→A1/E1；Agent→A2/A6；多模态/Science→A3/C2；Infra→A4/B4",
          "具身→A5；Builders→E2/A7；AGI/治理→E3/B2/F3；内容→D1–D5；出海→F1/F2",
          "正文27页版式与日历以贵方上传PDF为准",
          "附录补入原30场交付的执行卡、案例实拍、关联度口径与WAIC映射",
          "形成可签约、可执行、可验收的一套融合材料"]),
    ]
    for fn, title, lines in apps:
        im, d = new_img()
        d.text((80, 60), title, fill=(26, 26, 46), font=font(40))
        y = 180
        for line in lines:
            d.ellipse([90, y + 10, 110, y + 30], fill=(30, 77, 140))
            d.text((130, y), line[:70], fill=(26, 26, 46), font=font(26))
            y += 70
        d.text((80, 1020), "上海创智汇 × 同浦汇 · 融合版附录", fill=(138, 138, 154), font=font(18))
        im.save(os.path.join(outdir, fn))


def build_excel():
    wb = Workbook()
    thin = Border(
        left=Side(style="thin", color="D0D0D0"), right=Side(style="thin", color="D0D0D0"),
        top=Side(style="thin", color="D0D0D0"), bottom=Side(style="thin", color="D0D0D0"),
    )
    head_fill = PatternFill("solid", fgColor="1E4D8C")
    alt = PatternFill("solid", fgColor="EEF2F8")

    def hdr(ws, headers):
        for j, h in enumerate(headers, 1):
            c = ws.cell(1, j, h)
            c.font = Font(name="微软雅黑", bold=True, color="FFFFFF", size=10)
            c.fill = head_fill
            c.alignment = Alignment(wrap_text=True, vertical="center", horizontal="center")
            c.border = thin

    def rows_write(ws, data, widths, h=44):
        for i, row in enumerate(data, 2):
            for j, v in enumerate(row, 1):
                c = ws.cell(i, j, v)
                c.font = Font(name="微软雅黑", size=9)
                c.alignment = Alignment(wrap_text=True, vertical="center")
                c.border = thin
                if i % 2 == 0: c.fill = alt
            ws.row_dimensions[i].height = h
        ws.row_dimensions[1].height = 28
        for i, w in enumerate(widths, 1):
            ws.column_dimensions[get_column_letter(i)].width = w

    ws = wb.active; ws.title = "01-30场总表（PDF日历）"
    hdr(ws, ["编号", "主题", "活动名称", "档期", "形式", "人数", "负责人占比", "招商落点", "场地", "行程", "嘉宾", "价值", "交付"])
    data = [[e["code"], e["cat"], e["name"], e["month"], e["format"], e["people"], e["dm"], e["lease"],
             e["floor"], e["agenda"], e["guests"], e["value"], e["deliver"]] for e in EVENTS]
    rows_write(ws, data, [8, 6, 28, 16, 12, 10, 10, 22, 14, 40, 28, 20, 28], 52)

    ws2 = wb.create_sheet("02-价值交付分工收费")
    hdr(ws2, ["模块", "内容"])
    rows_write(ws2, [
        ["活动价值", "议题日常化；空间主题坐实；载体考核/补贴台账；看场转化；8间待租导流"],
        ["具体交付", "方案议程嘉宾函报名签到意向动线回访报价复盘照片通稿；月季年报"],
        ["我方主责", "策划邀约物料执行看场回访台账传播"],
        ["园区支持", "审定档期名单场地物业样板间价格审批入驻交割"],
        ["服中心/高校", "券务申报辅导；成果团队学者推荐；联合会渠道"],
        ["办公租金", "3.3 元/㎡/天"],
        ["物业费", "13.8 元/㎡/月"],
        ["待租房源", "8 间（活动优先看场）"],
        ["招商佣金", "成交首年有效净租金 × 2 个月"],
        ["活动服务费", "另行商定（年度打包或按场）"],
        ["关联度口径", "定向邀约 100% 主题关联；不设随机散客；负责人占比按场次目标"],
    ], [18, 70], 36)

    ws3 = wb.create_sheet("03-往期案例")
    hdr(ws3, ["类型", "案例名称", "时间地点", "机构", "企业", "议程", "照片", "映射"])
    rows_write(ws3, [
        [CASES[0]["title"], CASES[0]["name"], CASES[0]["meta"], CASES[0]["orgs"], CASES[0]["companies"],
         " | ".join(f"{a}:{b}" for a, b in CASES[0]["agenda"]), "；".join(CASES[0]["photos"]), CASES[0]["result"]],
        [CASES[1]["title"], CASES[1]["name"], CASES[1]["meta"], CASES[1]["orgs"], CASES[1]["companies"],
         " | ".join(f"{a}:{b}" for a, b in CASES[1]["agenda"]), "；".join(CASES[1]["photos"]), CASES[1]["result"]],
    ], [18, 36, 28, 36, 36, 50, 36, 36], 70)

    ws4 = wb.create_sheet("04-执行标准与转化")
    hdr(ws4, ["节点", "关键动作", "责任方", "交付物"])
    rows_write(ws4, [
        ["T-14", "锁定嘉宾议程物料报名", "同浦汇", "议程单、嘉宾确认函、报名页"],
        ["T-7", "邀约确认+场地安保报备", "同浦汇+园区", "到场名单、场地确认单"],
        ["T-1", "彩排+看场动线+礼包就位", "同浦汇+园区", "彩排记录、动线图"],
        ["活动日", "签到→主区→看场≥20分钟→洽谈", "全体", "签到表、意向登记表"],
        ["T+1~7", "三轮跟进，48小时报价", "同浦汇+服中心", "回访记录、报价单、复盘纪要"],
        ["转化指标", "留资≥60%；看场≥25%；月度看板", "双方联席", "转化月报"],
    ], [12, 36, 18, 28], 36)

    files = [
        "创智汇年度活动运营方案-融合版.xlsx",
        "chuangzhihui-annual-ops-plan-fused.xlsx",
        "创智汇30场活动方案排期表.xlsx",
        "chuangzhihui-30-events-schedule.xlsx",
        "创智汇30场活动方案-表格应答版.xlsx",
        "chuangzhihui-30-events-table-v3.xlsx",
    ]
    for fn in files:
        p = os.path.join(HERE, fn)
        wb.save(p)
        shutil.copy2(p, os.path.join(ART, fn))
    print("Excel OK")


if __name__ == "__main__":
    build_ppt()
    build_excel()
    build_pdf_from_ppt_pages()
    print("DONE")
