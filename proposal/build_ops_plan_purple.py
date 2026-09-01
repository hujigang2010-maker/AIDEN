# -*- coding: utf-8 -*-
"""创智汇年度活动运营方案 · 紫金主题版（双会资源并入）

视觉：沿用此前皇家紫金（深紫渐变 + 蓝紫点缀 + 香槟金）
内容：30场定义核 + WAIC/ChinaJoy招商向资源导入 + 修订口径
结构：背景→日历→六大主题→执行→转化收费案例→附录
"""
import os, shutil
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

HERE = os.path.dirname(os.path.abspath(__file__))
ART = "/opt/cursor/artifacts"
PHOTO = os.path.join(HERE, "cases", "photos")
FONT = "Microsoft YaHei"

BG_A, BG_B = "1A1140", "3A206E"
PANEL = RGBColor(0x27, 0x1A, 0x4D)
PANEL2 = RGBColor(0x1E, 0x15, 0x42)
INK = RGBColor(0xF5, 0xF1, 0xFF)
MUT = RGBColor(0xC6, 0xBA, 0xEA)
SOFT = RGBColor(0x94, 0x86, 0xC4)
ACC = RGBColor(0x8B, 0x7B, 0xFF)
ACC2 = RGBColor(0xC7, 0x7D, 0xFF)
GOLD = RGBColor(0xE9, 0xC2, 0x7C)
GREEN = RGBColor(0x56, 0xD6, 0xB6)
ROSE = RGBColor(0xF5, 0x8B, 0xB8)
LINE = RGBColor(0x4A, 0x3C, 0x7C)
ROW_A = RGBColor(0x2A, 0x1C, 0x52)
ROW_B = RGBColor(0x23, 0x17, 0x48)

SW, SH = Inches(13.333), Inches(7.5)
ML, CW = Inches(0.7), Inches(11.9)

EVENTS = [
    dict(code="P1", name="创智汇项目推广日①·中介与渠道专场", month="2026年8月上旬", fmt="项目推广日",
         people="≤30人", dm="≤30%", lease="中介渠道带看与项目推介", layer="L0",
         define="渠道点火场", src="园区自建", res="待租8间看场清单；双会转介话术",
         agenda="项目介绍→空间导览→政策要点→洽谈→建群", guests="中介/渠道/园区招商"),
    dict(code="P2", name="创智汇项目推广日②·五大行与金融机构专场", month="2026年8月中旬", fmt="项目推广日",
         people="≤30人", dm="≤30%", lease="金融资源链接与客户导入", layer="L0",
         define="金融转介场", src="园区自建+投研", res="WAIC金融议题摘要；CJ专业观众画像",
         agenda="项目路演→金融对接→座谈→看场→建群", guests="五大行及合作金融机构"),
    dict(code="P3", name="创智汇项目推广日③·投研机构专场", month="2026年8月中旬", fmt="项目推广日",
         people="≤30人", dm="≤30%", lease="投研视角放大项目影响力", layer="L0",
         define="叙事定调场", src="园区自建", res="WAIC赛道结论；CJ「与AI同游」判断",
         agenda="定位→赛道研判→圆桌→看场→建群", guests="券商研究所/产业研究院/智库"),
    dict(code="P4", name="创智汇项目推广日④·政府部门与载体协同专场", month="2026年8月下旬", fmt="项目推广日",
         people="≤30人", dm="≤30%", lease="政策与载体协同落地", layer="L0",
         define="政策协同场", src="云创+科企联/服中心", res="算力券/创新券；YOUNG立方政策包",
         agenda="载体介绍→政策协同→云创资源→看场→建群", guests="区部门/载体/云创/服中心"),
    dict(code="E1", name="WAIC成果承接·创智汇开放交流日", month="2026年8月下旬", fmt="开放交流日",
         people="≤30人", dm="≤30%", lease="衔接WAIC议题与园区看场", layer="L1",
         define="WAIC余热承接", src="WAIC总表", res="A/B线索抽样；ACT议题切片",
         agenda="议题速递→园区承接→双展台→洽谈→建群", guests="AI/内容企业、媒体、云创"),
    dict(code="A2", name="Agent智能体搭建一日营", month="2026年9月上旬", fmt="实训营",
         people="≤30人", dm="≤30%", lease="Agent初创落位", layer="L2",
         define="3F入孵核：Agent→OPC", src="WAIC·Agent", res="Agent池抽样；ACT-001缩尺",
         agenda="签到→工作流→实操→Demo→看场", guests="Agent工程师、初创团队"),
    dict(code="C1", name="ChinaJoy主题·数字娱乐与IP授权对接会", month="2026年9月中旬", fmt="对接会",
         people="≤30人", dm="≤30%", lease="CJ赛道企业/IP方入驻与展位", layer="L2",
         define="5F内容核：IP授权对接", src="CJ·游戏/IP", res="中小工作室+IP方；可触达263抽样",
         agenda="赛道解读→IP授权→撮合→看场→建群", guests="游戏/动漫/潮玩IP方、渠道"),
    dict(code="F1", name="出海专题活动（另议·另计价）", month="档期另议", fmt="专题（另议）",
         people="≤30人", dm="≤30%", lease="出海服务（不在年度包内）", layer="另计价",
         define="另议核：出海不进年度包", src="CJ国际展团（另案）", res="BTOB外资线索仅作另案素材",
         agenda="另案确定", guests="另案确定"),
    dict(code="A3", name="多模态智能体工作坊", month="2026年10月上旬", fmt="工作坊",
         people="≤30人", dm="≤30%", lease="多模态团队入驻", layer="L2",
         define="3F技术核：多模态看场", src="WAIC·多模态/AIGC", res="AIGC子类抽样；内容论坛切片",
         agenda="案例→实操→共创→评审→看场", guests="多模态应用团队、产品负责人"),
    dict(code="E4", name="ChinaJoy主题日·数字娱乐生态交流", month="2026年10月中旬", fmt="主题日",
         people="≤30人", dm="≤30%", lease="CJ生态企业集中看场", layer="L2",
         define="5F生态核：数字娱乐看场", src="CJ·游戏风云", res="Express/独立游戏中小团队池",
         agenda="开幕→赛道分享→展洽→看场→建群", guests="数字娱乐企业、内容工作室"),
    dict(code="B2", name="AI治理与可信智能体沙龙", month="2026年10月下旬", fmt="沙龙",
         people="≤30人", dm="≤30%", lease="合规型企业信任导入", layer="L2",
         define="信任核：治理合规", src="WAIC·治理论坛", res="治理议题摘要；合规顾问短名单",
         agenda="治理框架→合规清单→案例→对接", guests="合规顾问、企业法务/负责人"),
    dict(code="A4", name="火山引擎×算力Infra实务营", month="2026年11月上旬", fmt="厂商联训",
         people="≤30人", dm="≤30%", lease="高算力企业定向看场", layer="L2",
         define="算力核：云创×厂商联训", src="WAIC·算力", res="ACT-003缩尺；算力展商池抽样",
         agenda="政策包→代金券→案例→诊断→看场", guests="火山引擎、云创基地、技术负责人"),
    dict(code="B4", name="创新券·算力券·模型券实务沙龙", month="2026年11月中旬", fmt="实务沙龙",
         people="≤30人", dm="≤30%", lease="用券企业聚集", layer="L2",
         define="券务核：用券聚集转化", src="云创+WAIC算力", res="三券核销路径；云厂商联系人抽样",
         agenda="三券规则→核销→案例→开户→入驻洽谈", guests="券平台、云厂商、企业负责人"),
    dict(code="L1", name="领事到访接待（另计价）", month="档期另议", fmt="外事接待（另计价）",
         people="另议", dm="—", lease="领事到访（三部分费用之外）", layer="另计价",
         define="外事另计价核（到访≠挂牌）", src="另案", res="不占用年度30场资源池",
         agenda="另案确定", guests="领事及相关方（另案）"),
    dict(code="A5", name="具身智能空间交互体验日", month="2026年12月上旬", fmt="体验日",
         people="≤30人", dm="≤30%", lease="具身/机器人团队看场", layer="L3",
         define="双会叠加核：具身体验", src="WAIC具身+CJ Vision Future", res="ACT-002缩尺；场景观察+CJ前沿科技",
         agenda="演示→讲解→场景→踩点→洽谈", guests="具身团队、高校实验室"),
    dict(code="F4", name="创智汇AI年度Demo Day", month="2026年12月中旬", fmt="路演日",
         people="≤30人", dm="≤30%", lease="集中签约洽谈与媒体背书", layer="L4",
         define="年中收口核：精选签约", src="双会精选项目", res="年内导入A/B意向精选路演",
         agenda="开幕→精选路演→颁奖→洽谈", guests="投资人、链主、媒体、精选项目"),
    dict(code="D1", name="AIGC微短剧制片特训", month="2027年1月上旬", fmt="特训营",
         people="≤30人", dm="≤30%", lease="厂牌/工作室入驻", layer="L2",
         define="内容制片核", src="WAIC·AIGC", res="影视内容论坛切片；AIGC厂牌抽样",
         agenda="政策→制片→脚本→路演→看场", guests="导演、厂牌、平台方"),
    dict(code="A6", name="AI营销Agent实战营", month="2027年1月中旬", fmt="实战营",
         people="≤30人", dm="≤30%", lease="营销科技公司看场", layer="L2",
         define="营销科技核", src="WAIC·企业服务营销", res="企业服务/营销池抽样",
         agenda="策略→Agent→素材→复盘→转化", guests="投放操盘手、产品经理"),
    dict(code="B1", name="YOUNG立方×智能伙伴政策沙龙", month="2027年2月中旬", fmt="政策沙龙",
         people="≤30人", dm="≤30%", lease="内容/AI企业导入", layer="L2",
         define="政策获客核", src="服中心+双会客群", res="政策礼包+双会意向企业复邀",
         agenda="政策→礼包→画像→诊断→看场", guests="政策宣讲、服中心、企业负责人"),
    dict(code="B3", name="高企认定冲刺（AI企业专场）", month="2027年2月下旬", fmt="辅导会",
         people="≤30人", dm="≤30%", lease="待认定企业带政策入驻", layer="L2",
         define="带政策入驻核", src="服中心", res="双会AI应用类待认定抽样",
         agenda="条件→材料→初筛→套餐→激励", guests="辅导顾问、待认定企业负责人"),
    dict(code="A7", name="OPC超级个体黑客松（春）", month="2027年3月上旬", fmt="黑客松",
         people="≤30人", dm="≤30%", lease="获奖团队优先谈单元", layer="L2",
         define="OPC核：获奖谈单元", src="WAIC青年+CJ创作者", res="高校Builder；CJ AGS线索抽样",
         agenda="开题→开发→路演→礼包→看场", guests="评委、投资人、Builders"),
    dict(code="C2", name="高校成果转化·AI for Science日", month="2027年3月下旬", fmt="对接日",
         people="≤30人", dm="≤30%", lease="成果公司/实验室落户", layer="L2",
         define="成果转化核", src="学术支持+WAIC科研", res="复旦住房中心学术支持位；科研论坛切片",
         agenda="成果路演→场景→承接→看单元→洽谈", guests="复旦/同济技转、住房政策研究中心学者"),
    dict(code="F3", name="通往AGI季度圆桌", month="2027年3月下旬", fmt="闭门圆桌",
         people="≤30人", dm="≤30%", lease="研究型/模型团队", layer="L2",
         define="研究型客户核", src="WAIC·大模型论坛", res="ACT-001闭门形态；模型公司抽样",
         agenda="议题→圆桌→围炉→问答→看场", guests="学者、模型企业、投资人"),
    dict(code="D2", name="ChinaJoy主题·创作者内容首发①", month="2027年4月中旬", fmt="发布会",
         people="≤30人", dm="≤30%", lease="数字娱乐/内容品牌问询", layer="L2",
         define="创作者首发核", src="CJ·创作者展区", res="创作者/社团抽样；渠道玩法迁移",
         agenda="揭幕→发布→快闪→专访→招商通道", guests="创作者、媒体、渠道、CJ内容方"),
    dict(code="C3", name="ChinaJoy主题·游戏周边与潮玩供应链对接", month="2027年4月下旬", fmt="对接会",
         people="≤30人", dm="≤30%", lease="CJ供应链/周边企业展位落位", layer="L2",
         define="供应链核：周边/潮玩展位", src="CJ·魔玩/谷子", res="潮玩谷子高优先级；作品类示例",
         agenda="赛道介绍→供需对接→展位参观→报价→建群", guests="游戏周边、潮玩供应链、渠道采购"),
    dict(code="C4", name="专精特新·AI应用培育路演", month="2027年5月中旬", fmt="路演",
         people="≤30人", dm="≤30%", lease="成长型AI补位", layer="L2",
         define="成长补位核", src="WAIC·行业AI", res="ACT-006缩尺；行业AI池抽样",
         agenda="政策→路演→点评→金融→看场", guests="顾问、银行、基金、企业"),
    dict(code="D3", name="ChinaJoy主题·数字娱乐市集体验日", month="2027年5月下旬", fmt="体验日",
         people="≤30人", dm="≤30%", lease="优质摊主/工作室升级固定展位", layer="L2",
         define="市集转化核：摊主升级展位", src="CJ·创作者/市集", res="小微品牌；市集→固定展位路径",
         agenda="布展→体验→交流→转正洽谈→建群", guests="内容工作室、周边品牌、达人"),
    dict(code="C5", name="ChinaJoy主题·IP联名与衍生品撮合会", month="2027年6月中旬", fmt="撮合会",
         people="≤30人", dm="≤30%", lease="IP联名/衍生品团队入驻", layer="L2",
         define="IP商业化核", src="CJ·IP/联名", res="老字号×IP联名案例迁移；衍生品抽样",
         agenda="IP路演→联名模式→一对一→MOU→看场", guests="IP方、衍生品商、渠道"),
    dict(code="D4", name="ChinaJoy主题·沉浸式数字娱乐联展", month="2027年6月下旬", fmt="联展（短展期）",
         people="≤30人/场次", dm="≤30%", lease="闭幕洽谈集中转化", layer="L2",
         define="视觉落地核：短展期转化", src="CJ·氛围落地", res="联合内容方短名单；预约制控流",
         agenda="布展→预约观展→交流→闭幕洽谈", guests="联合内容方、渠道、品牌"),
    dict(code="D5", name="创作者首发②·衔接WAIC2027与ChinaJoy", month="2027年7月中旬", fmt="发布会",
         people="≤30人", dm="≤30%", lease="旺季补位，双大会预热", layer="L4",
         define="双会预热核", src="双会预热", res="全年双会线索复盘；预热下届双会",
         agenda="发布→渠道→看场→年框→双会预热", guests="IP方、渠道、媒体"),
]
assert len(EVENTS) == 30
EMAP = {e["code"]: e for e in EVENTS}

THEMES = [
    dict(key="A", name="智能体与算力训练", n="6场", role="3F基本盘：营/坊/联训/体验/黑客松",
         codes=["A2", "A3", "A4", "A5", "A6", "A7"], accent=ACC,
         blurb="对应WAIC：Agent·多模态·算力·具身；资源：ACT缩尺+线索抽样"),
    dict(key="B", name="政策与治理沙龙", n="4场", role="政策/券务/资质做成进园理由",
         codes=["B2", "B4", "B1", "B3"], accent=GOLD,
         blurb="对应WAIC治理/券务；服中心协同兑现"),
    dict(key="C", name="产业与ChinaJoy对接", n="5场", role="高校成果+CJ工作室/供应链/IP（合并潮玩类）",
         codes=["C1", "C2", "C3", "C4", "C5"], accent=ACC2,
         blurb="原东莞/汕头/扬州场次合并为ChinaJoy主题内容"),
    dict(key="D", name="AI内容与ChinaJoy人气", n="5场", role="5F人气引擎：特训/首发/市集/联展",
         codes=["D1", "D2", "D3", "D4", "D5"], accent=GREEN,
         blurb="对应CJ创作者经济+WAIC AIGC；与官方创作者活动错峰"),
    dict(key="E", name="启动月·项目推广与双会承接", n="6场", role="8月点火+WAIC承接+CJ主题日",
         codes=["P1", "P2", "P3", "P4", "E1", "E4"], accent=ROSE,
         blurb="原大场改为≤30人精准场；8月=项目推广日"),
    dict(key="F", name="收官·圆桌·另计价事项", n="2场+另计价", role="Demo Day收口；出海/领事另案",
         codes=["F4", "F3", "F1", "L1"], accent=GOLD,
         blurb="出海另议另计价；领事到访/挂牌另计价且分属不同性质"),
]

CASES = [
    {
        "sec": "05 · 往期案例（一）",
        "name": "杨「数」浦数字沙龙第七期：AI如何重塑企业DNA",
        "meta": "2025-06-30 · 美团上海综合指挥中心 · 政策/治理沙龙类",
        "orgs": "杨浦区委网信办；赛博院等",
        "companies": "美团；大众点评大模型团队等",
        "agenda": [("参观", "美团综合指挥中心数字化场景"), ("主题", "人机共生：AI如何重塑企业DNA"),
                   ("案例", "大众点评AI技术应用落地"), ("政策", "生成合成内容标识办法解读"), ("互动", "答疑与一对一交流")],
        "photos": ["case_c987bd91.jpg", "case_59157ea5.jpg", "case_6e51af9e.jpg"],
        "map": "对应：B类政策沙龙、A类训练营",
    },
    {
        "sec": "05 · 往期案例（二）",
        "name": "「融见科创·智启未来」人工智能专场路演暨投融资对接会",
        "meta": "2025-10 · 杨浦 · 邮储银行联合主办 · 路演对接类",
        "orgs": "杨浦科创促进会、邮储银行等",
        "companies": "复楚智能、卡房信息、一造科技、中科趋势、万笔千墨等",
        "agenda": [("开场", "促进会会长致辞与生态介绍"), ("主办", "邮储银行科技金融致辞"),
                   ("分享", "中邮证券AI行业趋势"), ("路演", "6+6+6+1对接模式"), ("收口", "区投促点评与招商邀请")],
        "photos": ["case_db1d5cac.jpg", "case_65931516.jpg", "case_cb0a3812.jpg"],
        "map": "对应：F4 Demo Day、项目推广日、路演类",
    },
]


def _grad(shape, stops, ang=90):
    spPr = shape._element.spPr
    for t in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill"):
        e = spPr.find(qn(t))
        if e is not None:
            spPr.remove(e)
    g = spPr.makeelement(qn("a:gradFill"), {})
    lst = g.makeelement(qn("a:gsLst"), {})
    for pos, col, al in stops:
        gs = g.makeelement(qn("a:gs"), {"pos": str(int(pos * 1000))})
        c = g.makeelement(qn("a:srgbClr"), {"val": col})
        if al is not None:
            c.append(g.makeelement(qn("a:alpha"), {"val": str(int(al * 1000))}))
        gs.append(c); lst.append(gs)
    g.append(lst)
    g.append(g.makeelement(qn("a:lin"), {"ang": str(int(ang * 60000)), "scaled": "1"}))
    ln = spPr.find(qn("a:ln"))
    (ln.addprevious(g) if ln is not None else spPr.append(g))


def slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.shadow.inherit = False; bg.line.fill.background()
    _grad(bg, [(0, BG_A, None), (55, BG_B, None), (100, BG_A, None)], 120)
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s


def rect(s, x, y, w, h, fill=None, line=None, lw=1.0, radius=False):
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    b.shadow.inherit = False
    if fill is None: b.fill.background()
    else: b.fill.solid(); b.fill.fore_color.rgb = fill
    if line is None: b.line.fill.background()
    else: b.line.color.rgb = line; b.line.width = Pt(lw)
    return b


def text(s, x, y, w, h, content, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = content; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = FONT
    return tb


def bullets(s, x, y, w, h, items, size=13, color=MUT, mark=ACC):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5); p.line_spacing = 1.15
        r0 = p.add_run(); r0.text = "▸ "; r0.font.size = Pt(size); r0.font.color.rgb = mark; r0.font.name = FONT
        r = p.add_run(); r.text = it; r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = FONT
    return tb


def header(s, sec, title, sub=None):
    text(s, ML, Inches(0.28), Inches(12), Inches(0.28), sec, size=13, color=ACC, bold=True)
    text(s, ML, Inches(0.55), Inches(12), Inches(0.42), title, size=24, color=INK, bold=True)
    rect(s, ML, Inches(1.05), Inches(0.55), Inches(0.04), fill=GOLD)
    if sub:
        text(s, ML, Inches(1.18), Inches(12), Inches(0.28), sub, size=12, color=MUT)


def footer(s, page, total):
    text(s, ML, Inches(7.1), Inches(9.5), Inches(0.25),
         "上海创智汇 × 同浦汇 · 年度活动运营方案 · 紫金版", size=11, color=SOFT)
    text(s, Inches(11.1), Inches(7.1), Inches(1.7), Inches(0.25),
         f"{page:02d} / {total:02d}", size=11, color=SOFT, align=PP_ALIGN.RIGHT)


def card(s, x, y, w, h, title, items, accent=ACC):
    rect(s, x, y, w, h, fill=PANEL, line=LINE, radius=True)
    rect(s, x, y, Inches(0.08), h, fill=accent)
    text(s, Emu(x + Inches(0.25)), Emu(y + Inches(0.15)), Emu(w - Inches(0.4)), Inches(0.32),
         title, size=15, color=INK, bold=True)
    bullets(s, Emu(x + Inches(0.25)), Emu(y + Inches(0.55)), Emu(w - Inches(0.45)), Emu(h - Inches(0.7)),
            items, size=13, mark=accent)


def table(s, y0, rows, col_w, sizes=None, rh=0.34):
    yy = y0
    if sizes is None: sizes = [11] * len(rows[0])
    for ri, row in enumerate(rows):
        xx = ML; h = Inches(0.4 if ri == 0 else rh)
        for ci, val in enumerate(row):
            cw = Emu(int(CW * col_w[ci]))
            if ri == 0:
                bg, fg, bd = PANEL, GOLD, True
            else:
                bg = ROW_A if ri % 2 == 0 else ROW_B
                fg, bd = INK, False
            rect(s, xx, yy, cw, h, fill=bg, line=LINE, lw=0.5)
            text(s, Emu(xx + Inches(0.05)), Emu(yy + Inches(0.05)), Emu(cw - Inches(0.08)), Emu(h - Inches(0.08)),
                 str(val), size=sizes[ci], color=fg, bold=bd, anchor=MSO_ANCHOR.MIDDLE)
            xx = Emu(xx + cw)
        yy = Emu(yy + h)
    return yy


def add_pic(s, path, x, y, w, h):
    if os.path.exists(path):
        s.shapes.add_picture(path, x, y, width=w, height=h)
    else:
        rect(s, x, y, w, h, fill=PANEL, line=LINE)


def mon(s):
    return s.replace("2026年", "").replace("2027年", "次")


def build():
    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
    TOTAL = 26
    n = 0

    def pg(s):
        nonlocal n
        n += 1
        footer(s, n, TOTAL)
        return n

    # 1 cover
    s = slide(prs)
    text(s, ML, Inches(1.6), Inches(12), Inches(0.3), "上海创智汇 × 同浦汇｜呈报园区", size=14, color=ACC, bold=True)
    text(s, ML, Inches(2.05), Inches(12), Inches(0.6), "创智汇年度活动运营方案", size=36, color=INK, bold=True)
    text(s, ML, Inches(2.75), Inches(12), Inches(0.35), "30场 · 定义核与双会资源导入 · 2026.08—2027.07", size=16, color=GOLD, bold=True)
    text(s, ML, Inches(3.25), Inches(12), Inches(0.35),
         "AI+数字内容无界共创港　｜　承接 WAIC · 衔接 ChinaJoy", size=14, color=MUT)
    card(s, ML, Inches(4.0), Inches(5.8), Inches(2.7), "本册要点", [
        "紫金主题视觉（沿用既往方案风格）",
        "每场≤30人 · 负责人≤30%",
        "同浦汇带客 / 园区销售",
        "WAIC→3F · ChinaJoy→5F",
        "主体：云创基地（国家级孵化器）",
    ], ACC)
    card(s, Inches(7.2), Inches(4.0), Inches(5.5), Inches(2.7), "边界", [
        "只取双会招商贴近切片，不堆全量名单",
        "出海另议另计价；领事到访/挂牌另计价",
        "不做租金租赁对赌与必要性要求",
        "不承诺必带外资企业",
        "学术支持：复旦住房政策研究中心",
    ], GOLD)
    pg(s)

    # 2 toc
    s = slide(prs); header(s, "目录", "本方案结构")
    table(s, Inches(1.45), [
        ["章节", "内容"],
        ["01 背景与总体思路", "为什么做、双会如何导入、做成什么样"],
        ["02 年度日历与运营节奏", "30场排期 + 四阶段 + 双会资源摘要"],
        ["03 六大主题活动详述", "A训练·B政策·C产业/CJ·D内容·E启动·F收官"],
        ["04 单场执行标准", "筹备→看场→建群→回访"],
        ["05 转化·分工·收费·案例", "带客/销售、付款、往期案例"],
        ["附录", "30场定义核与资源导入映射"],
    ], [0.28, 0.72], sizes=[14, 14], rh=0.52)
    pg(s)

    # 3 background
    s = slide(prs)
    header(s, "01 · 背景与判断", "大会闭幕之后，热度需要一个承接地",
           "主题响应：上海创智汇 · AI+数字内容无界共创港")
    card(s, ML, Inches(1.5), Inches(5.8), Inches(5.2), "判断", [
        "WAIC闭幕后：议题/嘉宾/线索需落到看场与洽谈",
        "ChinaJoy「与AI同游」：中小工作室/创作者/潮玩/具身可导入5F",
        "3F缺持续入孵节奏；5F缺持续内容人气",
        "因此：双会是资源库，30场是定义核",
        "每月两到三场，精而准，不断档",
    ], ACC)
    card(s, Inches(7.2), Inches(1.5), Inches(5.5), Inches(5.2), "运营目标", [
        "30场：全年活动总量",
        "12个月：2026.08—2027.07",
        "每场≤30人；全年约600人次+",
        "负责人占比≤30%",
        "不承诺必带外资企业",
    ], GOLD)
    pg(s)

    # 4 approach
    s = slide(prs); header(s, "01 · 总体思路", "一条主线：把双会的热度，变成园区的一年")
    card(s, ML, Inches(1.5), Inches(5.8), Inches(5.2), "园区得到什么", [
        "议题变日常：Agent/算力/治理/内容每月有场",
        "空间有主题：3F吃WAIC能力线，5F吃CJ内容线",
        "具身/AIGC为双会叠加带，跨楼层联动",
        "考核有支撑：签到、建群、回访摘要可归档",
        "云创基地国家级孵化器背书可感知",
    ], ACC)
    card(s, Inches(7.2), Inches(1.5), Inches(5.5), Inches(5.2), "招商得到什么", [
        "客群更准：按赛道抽样邀约，不堆全量名单",
        "动作标准化：每场固定看场+建群",
        "同浦汇带客，园区负责销售促成",
        "政策/券务活动中打包，降低决策门槛",
        "租金高于周边均价，故不做租赁对赌",
    ], GREEN)
    pg(s)

    # 5 framework
    s = slide(prs)
    header(s, "01 · 总体框架", "六大主题 + 四层叠加", "主题管内容；层次管资源导入顺序")
    table(s, Inches(1.45), [
        ["主题", "名称", "场次", "角色"],
        ["A", "智能体与算力训练", "6场", "3F基本盘"],
        ["B", "政策与治理沙龙", "4场", "政策/券务/资质进园理由"],
        ["C", "产业与ChinaJoy对接", "5场", "高校+CJ工作室/供应链/IP"],
        ["D", "AI内容与ChinaJoy人气", "5场", "5F人气：特训/首发/市集/联展"],
        ["E", "启动月·推广与双会承接", "6场", "项目推广日+WAIC+CJ主题日"],
        ["F", "收官·圆桌·另计价", "2场+另计价", "Demo Day；出海/领事另案"],
    ], [0.1, 0.32, 0.16, 0.42], sizes=[13, 13, 12, 13], rh=0.42)
    text(s, ML, Inches(5.6), CW, Inches(0.9),
         "资源叠加：L0点火(P1–P4) → L1承接(E1) → L2主线(A/B/C/D) → L3具身叠加(A5) → L4收口预热(F4/D5)",
         size=13, color=MUT)
    pg(s)

    # 6-7 tables
    for subset, title, sub in (
        (EVENTS[:15], "30场总表（上）", "2026年8月—12月　｜　每场≤30人 · 负责人≤30%"),
        (EVENTS[15:], "30场总表（下）", "2027年1月—7月　｜　每场≤30人 · 负责人≤30%"),
    ):
        s = slide(prs); header(s, "02 · 年度日历", title, sub)
        rows = [["编号", "活动名称", "档期", "形式", "人数", "层次", "招商落点"]]
        for e in subset:
            rows.append([e["code"], e["name"][:15], mon(e["month"])[:8], e["fmt"][:8],
                         e["people"][:6], e["layer"][:6], e["lease"][:11]])
        table(s, Inches(1.45), rows, [0.07, 0.27, 0.12, 0.12, 0.1, 0.1, 0.22], sizes=[10, 11, 10, 10, 10, 10, 10], rh=0.32)
        pg(s)

    # 8 rhythm + dual
    s = slide(prs); header(s, "02 · 运营节奏与双会资源", "四阶段重心 + 招商向资源摘要")
    card(s, ML, Inches(1.45), Inches(5.8), Inches(2.55), "四阶段", [
        "启动8–10月：推广日+WAIC+CJ开局",
        "攻坚11–12月：算力/券务/具身/Demo Day",
        "深耕1–3月：特训/政策/黑客松/成果",
        "收获4–7月：CJ内容密集+双会预热",
    ], ACC)
    card(s, Inches(7.2), Inches(1.45), Inches(5.5), Inches(2.55), "WAIC → 3F", [
        "963展商·4262品牌·3463联系人·A/B500",
        "ACT缩尺：闭门/具身/算力/行业AI",
        "用法：线索抽样，不堆全量",
    ], ACC2)
    card(s, ML, Inches(4.25), Inches(11.9), Inches(2.4), "ChinaJoy → 5F（主题「与AI同游」）", [
        "可触达263优先 · 最高：中小游戏工作室 · 高：Vision Future具身 / 潮玩谷子 / 创作者经济",
        "中腰部选址灵活；不硬追大厂总部；原潮玩产业集群场次合并为CJ主题内容",
    ], GOLD)
    pg(s)

    # 9-14 themes
    for t in THEMES:
        s = slide(prs)
        header(s, f"03 · 主题 {t['key']}", t["name"], t["blurb"])
        rows = [["编号", "活动", "档期", "形式", "招商落点", "资源导入"]]
        for c in t["codes"]:
            e = EMAP[c]
            rows.append([e["code"], e["name"][:13], mon(e["month"])[:8], e["fmt"][:7],
                         e["lease"][:11], e["res"][:16]])
        table(s, Inches(1.45), rows, [0.07, 0.24, 0.12, 0.12, 0.2, 0.25], sizes=[10, 11, 10, 10, 10, 10], rh=0.48)
        pg(s)

    # 15 execution
    s = slide(prs); header(s, "04 · 单场执行标准", "每场活动，同一套打法")
    table(s, Inches(1.45), [
        ["节点", "关键动作", "责任方", "交付物"],
        ["T-14", "锁定嘉宾议程物料 + 本场资源切片", "同浦汇", "议程/嘉宾确认/报名"],
        ["T-7", "邀约确认 + 场地报备", "同浦汇；园区场地确认", "名单/场地单"],
        ["T-1", "彩排 + 看场动线", "同浦汇；园区样板间", "彩排记录"],
        ["D日", "签到→主区→看场→建群", "同浦汇执行；园区支持", "签到/意向/群"],
        ["T+1~7", "带客回访", "同浦汇主责", "回访摘要知会园区"],
        ["销售促成", "房源条件与成交", "园区负责销售部分", "合同/交割"],
    ], [0.12, 0.34, 0.3, 0.24], sizes=[12, 13, 12, 12], rh=0.55)
    pg(s)

    # 16 convert + fee
    s = slide(prs); header(s, "05 · 转化 · 分工 · 收费", "带客/销售 · 付款唯一方案")
    card(s, ML, Inches(1.45), Inches(5.8), Inches(2.5), "转化与分工", [
        "同浦汇：策划、邀约、带客、建群、回访",
        "园区：场地物业、样板间、销售成交",
        "云创基地：主体运营与算力协同",
        "科企联/服中心：载体支持与政策手续",
        "D日建群→园区群内促成成交",
    ], GREEN)
    card(s, Inches(7.2), Inches(1.45), Inches(5.5), Inches(2.5), "收费口径", [
        "活动费：年度打包（线下确认）",
        "租金3.3元/㎡/天；物业13.8元/㎡/月",
        "建议免租期1–3个月（面议）",
        "不做租赁对赌与必要性要求",
        "出海/领事另计价",
    ], GOLD)
    table(s, Inches(4.2), [
        ["付款节点", "比例", "对应"],
        ["签约后7日内", "50%", "启动策划、年度日历锁定"],
        ["Q3 / Q4 / 次年Q1 / 次年Q2", "各10%（合计40%）", "与季度场次完成度挂钩"],
        ["年终收官后", "10%", "年报、台账、影像移交"],
    ], [0.34, 0.28, 0.38], sizes=[12, 12, 12], rh=0.42)
    pg(s)

    # 17-18 cases
    for c in CASES:
        s = slide(prs); header(s, c["sec"], c["name"], c["meta"])
        for j, fn in enumerate(c["photos"][:3]):
            add_pic(s, os.path.join(PHOTO, fn), Emu(ML + j * Inches(4.05)), Inches(1.45), Inches(3.9), Inches(2.15))
        rect(s, ML, Inches(3.8), Inches(7.3), Inches(2.9), fill=PANEL, line=LINE, radius=True)
        text(s, Emu(ML + Inches(0.2)), Inches(3.95), Inches(7.0), Inches(0.28), "当天议程", size=13, color=GOLD, bold=True)
        yy = Inches(4.35)
        for t0, d0 in c["agenda"]:
            text(s, Emu(ML + Inches(0.2)), yy, Inches(1.0), Inches(0.26), t0, size=11, color=ACC, bold=True)
            text(s, Emu(ML + Inches(1.2)), yy, Inches(5.8), Inches(0.26), d0, size=11, color=INK)
            yy = Emu(yy + Inches(0.26))
        rect(s, Inches(8.25), Inches(3.8), Inches(4.45), Inches(2.9), fill=PANEL2, line=LINE, radius=True)
        text(s, Inches(8.4), Inches(3.95), Inches(4.15), Inches(0.28), "可见机构/企业", size=13, color=GOLD, bold=True)
        text(s, Inches(8.4), Inches(4.35), Inches(4.15), Inches(2.15),
             f"{c['orgs']}\n\n企业侧：{c['companies']}\n\n映射：{c['map']}", size=12, color=MUT)
        pg(s)

    # 19 closing
    s = slide(prs); header(s, "结语", "把一年做满，双会热度落在园区")
    card(s, ML, Inches(1.5), Inches(5.8), Inches(5.2), "我们将交付", [
        "30场按日历执行（每场≤30人）",
        "8月项目推广日四场点火 + WAIC承接",
        "ChinaJoy主题内容场次（合并原潮玩类）",
        "每场定义核 + 资源导入标注",
        "带客、建群、回访摘要；园区做销售促成",
    ], ACC)
    card(s, Inches(7.2), Inches(1.5), Inches(5.5), Inches(5.2), "请园区协同", [
        "确认档期、场地物业与样板间",
        "共享在谈名单，群内发布条件促成成交",
        "云创基地/联合会/服中心协同到位",
        "政策礼包与券务按活动节奏提前到位",
        "月度联席会核对意向与下月排期",
    ], GOLD)
    pg(s)

    # 20 thanks
    s = slide(prs)
    text(s, ML, Inches(2.5), Inches(12), Inches(0.5), "谢谢审阅，期待与您共事", size=34, color=INK, bold=True, align=PP_ALIGN.CENTER)
    text(s, ML, Inches(3.3), Inches(12), Inches(0.4), "创智汇年度活动运营方案 · 紫金版", size=18, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    text(s, ML, Inches(3.9), Inches(12), Inches(0.35), "同浦汇带客 · 园区销售 · 衔接 WAIC & ChinaJoy", size=14, color=MUT, align=PP_ALIGN.CENTER)
    text(s, ML, Inches(4.6), Inches(12), Inches(1.0),
         "主体运营：上海市云计算创新基地（国家级孵化器）\n学术支持：复旦大学住房政策研究中心\n载体支持：杨浦区科技企业联合会 · 科技企业服务中心",
         size=13, color=SOFT, align=PP_ALIGN.CENTER)
    pg(s)

    # 21-23 appendix definition
    for subset, title in ((EVENTS[:10], "附录 · 定义核与资源导入（1/3）"),
                          (EVENTS[10:20], "附录 · 定义核与资源导入（2/3）"),
                          (EVENTS[20:], "附录 · 定义核与资源导入（3/3）")):
        s = slide(prs); header(s, "附录", title, "每场一句话定义 + 导入什么资源（招商向）")
        rows = [["编号", "定义核", "资源源", "资源导入"]]
        for e in subset:
            rows.append([e["code"], e["define"][:22], e["src"][:14], e["res"][:26]])
        table(s, Inches(1.45), rows, [0.08, 0.32, 0.2, 0.4], sizes=[11, 12, 11, 11], rh=0.42)
        pg(s)

    # 24-25 exec elements
    for subset, title in ((EVENTS[:15], "附录 · 逐场执行要素（上）"), (EVENTS[15:], "附录 · 逐场执行要素（下）")):
        s = slide(prs); header(s, "附录", title)
        rows = [["编号", "活动", "行程", "嘉宾"]]
        for e in subset:
            rows.append([e["code"], e["name"][:14], e["agenda"][:26], e["guests"][:16]])
        table(s, Inches(1.45), rows, [0.08, 0.26, 0.38, 0.28], sizes=[10, 11, 11, 11], rh=0.32)
        pg(s)

    # 26口径
    s = slide(prs); header(s, "附录 · 主体与口径一页纸", "合并后统一口径")
    table(s, Inches(1.45), [
        ["模块", "口径"],
        ["主体运营范围", "上海市云计算创新基地（国家级孵化器）"],
        ["学术支持", "复旦大学住房政策研究中心"],
        ["载体支持", "杨浦区科技企业联合会、科技企业服务中心"],
        ["执行带客 / 销售", "同浦汇带客；园区负责销售部分与合同交割"],
        ["人数 / 负责人", "每场≤30人；负责人≤30%；全年约600人次+"],
        ["双会原则", "WAIC→3F；ChinaJoy→5F；具身/AIGC叠加；只取招商切片"],
        ["外企 / 出海 / 领事", "不承诺必带外企；出海另议；领事到访与挂牌均另计价且分性质"],
        ["租金与付款", "不对赌；免租期1–3个月面议；付款50%+4×10%+10%"],
        ["视觉", "皇家紫金主题（深紫渐变 + 蓝紫点缀 + 香槟金）"],
    ], [0.28, 0.72], sizes=[12, 13], rh=0.4)
    pg(s)

    assert n == TOTAL, f"{n} != {TOTAL}"

    os.makedirs(ART, exist_ok=True)
    outs = [
        "创智汇年度活动运营方案-双会并入融合版.pptx",
        "chuangzhihui-annual-ops-plan-dual-merge.pptx",
        "创智汇年度活动运营方案-融合版.pptx",
        "chuangzhihui-annual-ops-plan-fused.pptx",
        "创智汇30场活动专项交付-双会资源融合版.pptx",
        "chuangzhihui-30-events-waic-cj-fuse.pptx",
    ]
    tmp = os.path.join(HERE, "_tmp_purple.pptx")
    prs.save(tmp)
    for fn in outs:
        p = os.path.join(HERE, fn)
        shutil.copy2(tmp, p)
        shutil.copy2(p, os.path.join(ART, fn))
    os.remove(tmp)
    print("PPT slides:", len(prs.slides))

    # Excel
    wb = Workbook()
    thin = Border(left=Side(style="thin", color="4A3C7C"), right=Side(style="thin", color="4A3C7C"),
                  top=Side(style="thin", color="4A3C7C"), bottom=Side(style="thin", color="4A3C7C"))
    head = PatternFill("solid", fgColor="3A206E")
    alt = PatternFill("solid", fgColor="F3F0FF")

    def hdr(ws, headers):
        for j, h in enumerate(headers, 1):
            c = ws.cell(1, j, h); c.font = Font(name="微软雅黑", bold=True, color="FFFFFF", size=10)
            c.fill = head; c.alignment = Alignment(wrap_text=True, vertical="center"); c.border = thin

    def write(ws, data, widths, h=44):
        for i, row in enumerate(data, 2):
            for j, v in enumerate(row, 1):
                c = ws.cell(i, j, v); c.font = Font(name="微软雅黑", size=9)
                c.alignment = Alignment(wrap_text=True, vertical="center"); c.border = thin
                if i % 2 == 0: c.fill = alt
            ws.row_dimensions[i].height = h
        ws.row_dimensions[1].height = 28
        for i, w in enumerate(widths, 1):
            ws.column_dimensions[get_column_letter(i)].width = w

    ws = wb.active; ws.title = "01-活动概况"
    hdr(ws, ["编号", "活动名称", "时间", "形式", "人数", "负责人", "层次", "招商落点", "行程", "嘉宾"])
    write(ws, [[e["code"], e["name"], e["month"], e["fmt"], e["people"], e["dm"], e["layer"],
                e["lease"], e["agenda"], e["guests"]] for e in EVENTS],
          [8, 32, 16, 12, 10, 8, 8, 22, 36, 28], 48)

    ws2 = wb.create_sheet("02-定义核资源导入")
    hdr(ws2, ["编号", "活动名称", "层次", "定义核", "资源源", "资源导入"])
    write(ws2, [[e["code"], e["name"], e["layer"], e["define"], e["src"], e["res"]] for e in EVENTS],
          [8, 32, 10, 28, 22, 40], 48)

    ws3 = wb.create_sheet("03-六大主题")
    hdr(ws3, ["主题", "名称", "场次编号", "说明"])
    write(ws3, [[t["key"], t["name"], "、".join(t["codes"]), t["blurb"]] for t in THEMES],
          [8, 28, 28, 50], 40)

    ws4 = wb.create_sheet("04-主体口径")
    hdr(ws4, ["模块", "口径"])
    write(ws4, [
        ["视觉", "皇家紫金主题（深紫渐变+蓝紫点缀+香槟金）"],
        ["主体运营范围", "上海市云计算创新基地（国家级孵化器）"],
        ["学术支持", "复旦大学住房政策研究中心"],
        ["载体支持", "杨浦区科技企业联合会、科技企业服务中心"],
        ["执行带客", "同浦汇"],
        ["销售", "园区负责销售部分"],
        ["人数", "每场≤30人；全年约600人次+"],
        ["负责人", "≤30%"],
        ["双会", "WAIC→3F；ChinaJoy→5F；具身/AIGC叠加"],
        ["出海/领事", "另议另计价；到访≠挂牌"],
        ["租金", "3.3元/㎡/天；不做对赌；免租期1–3个月面议"],
        ["付款", "签约后7日50%；四季各10%；年终10%"],
    ], [18, 72], 34)

    for fn in ["创智汇年度活动运营方案-双会并入融合版.xlsx", "chuangzhihui-annual-ops-plan-dual-merge.xlsx",
               "创智汇年度活动运营方案-融合版.xlsx", "chuangzhihui-annual-ops-plan-fused.xlsx",
               "创智汇30场活动专项交付-双会资源融合版.xlsx", "chuangzhihui-30-events-waic-cj-fuse.xlsx"]:
        p = os.path.join(HERE, fn); wb.save(p); shutil.copy2(p, os.path.join(ART, fn))
    print("Excel OK")


if __name__ == "__main__":
    build()
