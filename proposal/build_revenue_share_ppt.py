# -*- coding: utf-8 -*-
"""创智汇 · 科创科技服务中心收益评估与分成方案 PPT
甲方：同普会　乙方：杨浦区科技企业服务中心
用法: python3 build_revenue_share_ppt.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
ART = "/opt/cursor/artifacts"
FONT = "Microsoft YaHei"
FONT_EN = "Arial"
BG_A, BG_B = "1A1140", "3A206E"
INK = RGBColor(0xF5, 0xF1, 0xFF)
MUT = RGBColor(0xC6, 0xBA, 0xEA)
SOFT = RGBColor(0x94, 0x86, 0xC4)
ACC = RGBColor(0x8B, 0x7B, 0xFF)
ACC2 = RGBColor(0xC7, 0x7D, 0xFF)
GOLD = RGBColor(0xE9, 0xC2, 0x7C)
GREEN = RGBColor(0x56, 0xD6, 0xB6)
ROSE = RGBColor(0xF5, 0x8B, 0xB8)
LINE = RGBColor(0x4A, 0x3C, 0x7C)
SW, SH = Inches(13.333), Inches(7.5)
ML, CW = Inches(0.85), Inches(11.63)
prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]
_FOOT = []
IDX = 0


def _grad(shape, stops, ang=90):
    spPr = shape._element.spPr
    for t in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill"):
        e = spPr.find(qn(t))
        if e is not None:
            spPr.remove(e)
    g = spPr.makeelement(qn("a:gradFill"), {})
    lst = g.makeelement(qn("a:gsLst"), {})
    for pos, col, al in stops:
        gs = g.makeelement(qn("a:gs"), {"pos": str(int(pos * 1000))})
        c = g.makeelement(qn("a:srgbClr"), {"val": col})
        if al is not None:
            c.append(g.makeelement(qn("a:alpha"), {"val": str(int(al * 1000))}))
        gs.append(c)
        lst.append(gs)
    g.append(lst)
    g.append(g.makeelement(qn("a:lin"), {"ang": str(int(ang * 60000)), "scaled": "1"}))
    ln = spPr.find(qn("a:ln"))
    (ln.addprevious(g) if ln is not None else spPr.append(g))


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.shadow.inherit = False
    bg.line.fill.background()
    _grad(bg, [(0, BG_A, None), (55, BG_B, None), (100, BG_A, None)], 120)
    s.shapes._spTree.remove(bg._element)
    s.shapes._spTree.insert(2, bg._element)
    return s


def rect(s, x, y, w, h, fill=None, line=None, lw=1.0, radius=False, grad=None, gang=90):
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    b.shadow.inherit = False
    if grad is not None:
        _grad(b, grad, gang)
    elif fill is None:
        b.fill.background()
    else:
        b.fill.solid()
        b.fill.fore_color.rgb = fill
    if line is None:
        b.line.fill.background()
    else:
        b.line.color.rgb = line
        b.line.width = Pt(lw)
    return b


def text(s, x, y, w, h, runs, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, space=1.0, font=FONT, spacing=None):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    if isinstance(runs, str):
        runs = [(runs, color, bold)]
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = space
    for seg in runs:
        t, c, b = seg[0], seg[1], seg[2]
        f = seg[3] if len(seg) > 3 else font
        r = p.add_run()
        r.text = t
        r.font.size = Pt(size)
        r.font.bold = b
        r.font.color.rgb = c
        r.font.name = f
        if spacing is not None:
            r._r.get_or_add_rPr().set("spc", str(int(spacing * 100)))
    return tb


def bullets(s, x, y, w, h, items, size=14, color=MUT, gap=7, mark=ACC, lh=1.18):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_top = Pt(2)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = lh
        r0 = p.add_run()
        r0.text = "▪  "
        r0.font.size = Pt(size)
        r0.font.color.rgb = mark
        r0.font.name = FONT
        r = p.add_run()
        r.text = it
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = FONT
    return tb


def header(s, sec, eyebrow, title):
    rect(s, ML, Inches(0.58), Inches(0.06), Inches(0.82), fill=GOLD)
    text(s, Emu(ML + Inches(0.22)), Inches(0.56), Inches(9), Inches(0.32),
         eyebrow, size=12, color=GOLD, bold=True, spacing=2, font=FONT_EN)
    text(s, Emu(ML + Inches(0.22)), Inches(0.86), Inches(10.5), Inches(0.58),
         title, size=24, color=INK, bold=True)
    text(s, Inches(11.35), Inches(0.48), Inches(1.5), Inches(0.85),
         [(sec, GOLD, True, FONT_EN)], size=36, align=PP_ALIGN.RIGHT, font=FONT_EN)
    rect(s, Inches(11.55), Inches(1.26), Inches(1.28), Pt(2), fill=LINE)
    rect(s, ML, Inches(6.92), CW, Pt(1), fill=LINE)


def footer(s, idx):
    text(s, ML, Inches(7.0), Inches(9.2), Inches(0.32),
         [("上海创智汇 ", SOFT, False), ("CHUANGZHIHUI", SOFT, False, FONT_EN),
          ("  ·  科创服务中心收益评估与分成方案", SOFT, False)], size=9.5)
    rs = text(s, Inches(11.15), Inches(7.0), Inches(1.7), Inches(0.32),
              [("%02d" % idx, GOLD, True, FONT_EN), (" / 12", SOFT, False, FONT_EN)],
              size=10.5, align=PP_ALIGN.RIGHT)
    _FOOT.append(rs.text_frame.paragraphs[0].runs[1])


def card(s, x, y, w, h, title=None, body=None, items=None, accent=ACC, tsize=14.5, bsize=12):
    rect(s, x, y, w, h, grad=[(0, "2A1E55", None), (100, "1E1542", None)], gang=120, line=LINE, lw=1, radius=True)
    rect(s, x, y, Inches(0.07), h, fill=accent)
    if title:
        text(s, Emu(x + Inches(0.28)), Emu(y + Inches(0.16)), Emu(w - Inches(0.45)), Inches(0.42),
             title, size=tsize, color=INK, bold=True)
    if body:
        text(s, Emu(x + Inches(0.28)), Emu(y + Inches(0.58)), Emu(w - Inches(0.5)), Emu(h - Inches(0.7)),
             body, size=bsize, color=MUT, space=1.15)
    if items:
        bullets(s, Emu(x + Inches(0.28)), Emu(y + Inches(0.55)), Emu(w - Inches(0.5)), Emu(h - Inches(0.65)),
                items, size=bsize, mark=accent)


def table(s, x, y, w, rows, col_w, sizes=None, rh=Inches(0.42), head_rh=Inches(0.46)):
    n = len(rows[0])
    if sizes is None:
        sizes = [11] * n
    yy = y
    for ri, row in enumerate(rows):
        xx = x
        h = head_rh if ri == 0 else rh
        for ci, val in enumerate(row):
            cw = Inches(col_w[ci]) if not isinstance(col_w[ci], type(Inches(1))) else col_w[ci]
            # col_w as fractions of w
            if isinstance(col_w[ci], float):
                cw = Emu(int(w * col_w[ci]))
            cell = rect(s, xx, yy, cw, h,
                        fill=RGBColor(0x2A, 0x1E, 0x55) if ri == 0 else (RGBColor(0x22, 0x18, 0x48) if ri % 2 else RGBColor(0x1C, 0x14, 0x3A)),
                        line=LINE, lw=0.75)
            c = GOLD if ri == 0 else INK
            b = True if ri == 0 or ci == 0 else False
            text(s, Emu(xx + Inches(0.08)), Emu(yy + Inches(0.06)), Emu(cw - Inches(0.12)), Emu(h - Inches(0.08)),
                 str(val), size=sizes[ci], color=c, bold=b, anchor=MSO_ANCHOR.MIDDLE)
            xx = Emu(xx + cw)
        yy = Emu(yy + h)
    return yy


def nxt():
    global IDX
    IDX += 1
    return IDX


# ========== 01 封面 ==========
s = slide()
ov = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.2), Inches(-1.5), Inches(5), Inches(5))
ov.shadow.inherit = False
ov.line.fill.background()
_grad(ov, [(0, "5B3FA8", 40), (100, "1A1140", 0)])
text(s, ML, Inches(1.6), Inches(11), Inches(0.4),
     [("REVENUE SHARE  ·  SERVICE CENTER", GOLD, True, FONT_EN)], size=14, spacing=3)
text(s, ML, Inches(2.15), Inches(11.5), Inches(1.1),
     "科创科技服务中心 · 收益评估与分成方案", size=34, color=INK, bold=True)
text(s, ML, Inches(3.3), Inches(11), Inches(0.45),
     "政策申报收益项测算  ×  门类分成方案  ×  贡献比例分配", size=16, color=MUT)
rect(s, ML, Inches(4.0), Inches(2.2), Pt(3), fill=GOLD)

# parties
rect(s, ML, Inches(4.4), Inches(5.4), Inches(1.55),
     grad=[(0, "2A1E55", None), (100, "1E1542", None)], line=LINE, lw=1, radius=True)
text(s, Emu(ML + Inches(0.3)), Inches(4.55), Inches(5), Inches(0.3),
     [("甲方", GOLD, True)], size=12)
text(s, Emu(ML + Inches(0.3)), Inches(4.9), Inches(5), Inches(0.7),
     "同普会\n（协同同浦汇产业运营资源）", size=16, color=INK, bold=True, space=1.2)

rect(s, Inches(7.0), Inches(4.4), Inches(5.4), Inches(1.55),
     grad=[(0, "2A1E55", None), (100, "1E1542", None)], line=LINE, lw=1, radius=True)
text(s, Inches(7.3), Inches(4.55), Inches(5), Inches(0.3),
     [("乙方（我们）", GOLD, True)], size=12)
text(s, Inches(7.3), Inches(4.9), Inches(5), Inches(0.7),
     "杨浦区科技企业服务中心\n（科企联体系 · 政务通道与申报主体）", size=15, color=INK, bold=True, space=1.2)

text(s, ML, Inches(6.3), Inches(11), Inches(0.4),
     "上海创智汇 · AI+数字内容无界共创港　｜　沟通版 v1.0　｜　测算供协商，最终以协议为准",
     size=12, color=SOFT)

# ========== 02 目录 ==========
idx = nxt()
s = slide()
header(s, "00", "AGENDA", "汇报目录 · 评估逻辑")
footer(s, idx)
items = [
    ("01", "合作双方与评估边界", "谁出什么、评什么、钱从哪来"),
    ("02", "政策申报收益项清单", "载体 / 活动 / 奖项 / 其他政策"),
    ("03", "服务方可实现收益评估", "可申请金额 · 服务费比例/金额"),
    ("04", "收益分成方案", "按门类设定甲乙分成比例"),
    ("05", "贡献度与测算结果", "成熟年 / 三年累计分配"),
    ("06", "结论与签约建议", "分成原则写入补充协议"),
]
for i, (num, t, sub) in enumerate(items):
    y = Inches(1.55) + Inches(i * 0.82)
    rect(s, ML, y, CW, Inches(0.72),
         grad=[(0, "2A1E55", None), (100, "1E1542", None)], line=LINE, lw=1, radius=True)
    text(s, Emu(ML + Inches(0.3)), Emu(y + Inches(0.12)), Inches(0.8), Inches(0.5),
         [(num, GOLD, True, FONT_EN)], size=22)
    text(s, Emu(ML + Inches(1.2)), Emu(y + Inches(0.1)), Inches(6), Inches(0.35),
         t, size=16, color=INK, bold=True)
    text(s, Emu(ML + Inches(1.2)), Emu(y + Inches(0.4)), Inches(9), Inches(0.28),
         sub, size=12, color=MUT)

# ========== 03 双方与边界 ==========
idx = nxt()
s = slide()
header(s, "01", "PARTIES & SCOPE", "合作双方与评估边界")
footer(s, idx)
card(s, ML, Inches(1.5), Inches(5.7), Inches(2.4),
     title="甲方 · 同普会",
     items=[
         "产业资源、项目导入、商务统筹",
         "同浦汇协同：招商/出海/活动运营",
         "客户关系与品牌传播协同",
         "贡献侧重：获客 · 资源 · 商务转化",
     ], accent=GOLD)
card(s, Inches(6.9), Inches(1.5), Inches(5.55), Inches(2.4),
     title="乙方 · 杨浦区科技企业服务中心",
     items=[
         "政务通道、载体/服务机构资质",
         "政策申报执行、材料与过审",
         "创新券平台核销与合规把关",
         "贡献侧重：资质 · 申报 · 兑现落地",
     ], accent=ACC)
card(s, ML, Inches(4.1), Inches(11.6), Inches(2.5),
     title="本次评估边界（服务中心未来收益）",
     items=[
         "评估对象：挂牌「科技企业服务中心」后，因政策申报与载体运营可沉淀到服务方的现金收益",
         "三块主账：①载体可申请/到账　②活动经费可申请/到账　③申报奖项带来的服务费+成功分成",
         "其他政策项：房租补贴、创新券、政府专项等——按服务费比例或到账分成计入服务方收益",
         "不含：业主租金、招商佣金主合同（可另表）；本 PPT 聚焦服务中心政策服务收益与甲乙分佣",
     ], accent=GREEN, bsize=12.5)

# ========== 04 收益项总览 ==========
idx = nxt()
s = slide()
header(s, "02", "REVENUE MAP", "政策申报收益项总览（服务方可计）")
footer(s, idx)
rows = [
    ["门类", "详细项", "对企业/载体的政策金额", "服务方可实现收益口径"],
    ["载体", "成果转化平台运营补贴", "≤100万/年", "直接到平台 10–20万/年（认定+满年）"],
    ["载体", "创业孵化基地运营补贴", "10万/年（评估达标）", "直接到平台约10万/年"],
    ["活动", "YOUNG立方活动经费补贴", "投入×50%，封顶≤200万/年", "以申报主体到账；或执行服务费"],
    ["奖项", "高企 / 专精特新 / 小巨人", "20 / 10·30 / 50–100万（企业）", "服务费2–8万/家 + 到账5%–15%"],
    ["其他", "科小备案 / 创新券 / 房租补贴等", "资格或≤30–100万/年（企业）", "服务费0.3–2万 或 到账5%–15%"],
]
# manual compact table
yy = Inches(1.48)
cols = [0.10, 0.26, 0.30, 0.34]
for ri, row in enumerate(rows):
    xx = ML
    h = Inches(0.72) if ri == 0 else Inches(0.78)
    for ci, val in enumerate(row):
        cw = Emu(int(CW * cols[ci]))
        bgc = RGBColor(0x2A, 0x1E, 0x55) if ri == 0 else (RGBColor(0x24, 0x1A, 0x4A) if ri % 2 else RGBColor(0x1C, 0x14, 0x3A))
        rect(s, xx, yy, cw, h, fill=bgc, line=LINE, lw=0.7)
        text(s, Emu(xx + Inches(0.1)), Emu(yy + Inches(0.12)), Emu(cw - Inches(0.18)), Emu(h - Inches(0.18)),
             str(val), size=11 if ri else 12, color=GOLD if ri == 0 else INK, bold=(ri == 0 or ci == 0),
             anchor=MSO_ANCHOR.MIDDLE)
        xx = Emu(xx + cw)
    yy = Emu(yy + h)
text(s, ML, Inches(6.45), CW, Inches(0.35),
     "＊产业补贴主体多为企业/载体；服务中心收益＝直接运营补贴 + 申报服务费 + 成功分成（及活动主体到账）。",
     size=11, color=SOFT)

# ========== 05 详细收益评估 ==========
idx = nxt()
s = slide()
header(s, "02", "ITEMIZED ASSESSMENT", "分项收益评估 · 可申请 / 服务费 / 确定性")
footer(s, idx)
rows = [
    ["详细项", "可申请/政策金额", "服务费比例或金额", "成熟年服务方收益（示意）", "确定性"],
    ["载体运营补贴", "≤100万/年", "直接到账（非服务费）", "15–20万", "较确定·需认定"],
    ["孵化基地运营", "10万/年", "直接到账", "8–10万", "较确定·评估达标"],
    ["活动经费补贴", "≤200万/年", "主体到账或执行分成", "30–60万*", "依投入与主体"],
    ["高企认定", "企业20万/家", "服务费2–8万/家", "5万×8家≈40万", "确定（服务费）"],
    ["专精特新/小巨人", "10/30/50/100万", "服务费2–8万 + 分成5–15%", "服务费+分成≈45–80万", "服务费确定·分成依到账"],
    ["科小/创新券等", "资格或≤30万/年", "0.3–2万 或核销分成", "10–25万", "需入驻平台"],
    ["房租/专项补贴", "≤100万/年等", "到账成功费5%–15%", "15–40万", "依兑现"],
]
yy = Inches(1.42)
cols = [0.18, 0.18, 0.22, 0.26, 0.16]
for ri, row in enumerate(rows):
    xx = ML
    h = Inches(0.58) if ri == 0 else Inches(0.62)
    for ci, val in enumerate(row):
        cw = Emu(int(CW * cols[ci]))
        bgc = RGBColor(0x2A, 0x1E, 0x55) if ri == 0 else (RGBColor(0x24, 0x1A, 0x4A) if ri % 2 else RGBColor(0x1C, 0x14, 0x3A))
        rect(s, xx, yy, cw, h, fill=bgc, line=LINE, lw=0.6)
        text(s, Emu(xx + Inches(0.08)), Emu(yy + Inches(0.08)), Emu(cw - Inches(0.12)), Emu(h - Inches(0.12)),
             str(val), size=10.5 if ri else 11, color=GOLD if ri == 0 else INK,
             bold=(ri == 0 or ci == 0), anchor=MSO_ANCHOR.MIDDLE)
        xx = Emu(xx + cw)
    yy = Emu(yy + h)
text(s, ML, Inches(6.5), CW, Inches(0.3),
     "＊活动经费：若服务中心为申报主体且活动投入抬升，可申空间大；若仅基线30万投入，可申约15万。上表成熟年取中位示意。",
     size=10.5, color=SOFT)

# ========== 06 成熟年汇总 ==========
idx = nxt()
s = slide()
header(s, "03", "ANNUAL OUTLOOK", "成熟年服务方收益总览（合理中位）")
footer(s, idx)
kpis = [
    ("载体类", "25–30万", "运营补贴到平台", GOLD),
    ("活动类", "30–60万", "经费到账/执行分成", ACC),
    ("奖项服务", "85–120万", "服务费+成功分成", ACC2),
    ("其他政策", "25–65万", "创新券/房租/专项", GREEN),
]
for i, (t, n, sub, col) in enumerate(kpis):
    x = ML + Inches(i * 3.0)
    rect(s, x, Inches(1.55), Inches(2.85), Inches(2.15),
         grad=[(0, "2A1E55", None), (100, "1E1542", None)], line=LINE, lw=1, radius=True)
    text(s, Emu(x + Inches(0.2)), Inches(1.7), Inches(2.5), Inches(0.35), t, size=13, color=col, bold=True)
    text(s, Emu(x + Inches(0.2)), Inches(2.2), Inches(2.5), Inches(0.7), n, size=26, color=INK, bold=True)
    text(s, Emu(x + Inches(0.2)), Inches(3.0), Inches(2.5), Inches(0.5), sub, size=12, color=MUT)

rect(s, ML, Inches(3.95), CW, Inches(2.65),
     grad=[(0, "2A1E55", None), (100, "1E1542", None)], line=GOLD, lw=1.25, radius=True)
text(s, Emu(ML + Inches(0.4)), Inches(4.15), Inches(10.5), Inches(0.4),
     "成熟年服务方合计（中位）", size=14, color=GOLD, bold=True)
text(s, Emu(ML + Inches(0.4)), Inches(4.6), Inches(10.5), Inches(0.7),
     "约 165 – 275 万元 / 年", size=32, color=INK, bold=True)
bullets(s, Emu(ML + Inches(0.4)), Inches(5.4), Inches(10.8), Inches(1.0),
        ["保守：约165万（载体25+活动30+奖项85+其他25）　｜　积极：约275万（载体30+活动60+奖项120+其他65）",
         "前置条件：载体/成果转化平台认定 · 服务中心挂牌 · 创新券服务机构入驻 · 样板企业申报流水跑通",
         "以下分成方案均在「服务方已实现到账收益」基础上，再按门类在甲乙之间分配"],
        size=12.5, mark=GOLD, gap=5)

# ========== 07 分成原则 ==========
idx = nxt()
s = slide()
header(s, "04", "SHARE PRINCIPLES", "分成原则 · 按贡献定比例")
footer(s, idx)
card(s, ML, Inches(1.5), Inches(3.7), Inches(5.1),
     title="原则一 · 资质归谁、基础归谁",
     items=[
         "载体认定、服务中心挂牌、创新券机构资格主要由乙方持有",
         "直接到平台的载体运营补贴：乙方基础份额更高",
         "甲方以资源与商务换取对应分成，而非替代资质主体",
     ], accent=GOLD, bsize=13)
card(s, Inches(5.5), Inches(1.5), Inches(3.7), Inches(5.1),
     title="原则二 · 谁获客、谁执行",
     items=[
         "甲方/同浦汇：获客、活动引流、企业转化",
         "乙方：申报执行、材料、过审、合规",
         "服务费类：执行方（乙方）为主；导流方（甲方）按贡献取成",
     ], accent=ACC, bsize=13)
card(s, Inches(9.3), Inches(1.5), Inches(3.15), Inches(5.1),
     title="原则三 · 风险共担",
     items=[
         "成功分成随到账结算，未兑现不计",
         "季度对账、先扣直接成本再分",
         "重大专项可单独立账",
         "比例写入补充协议可调",
     ], accent=GREEN, bsize=13)

# ========== 08 门类分成方案 ==========
idx = nxt()
s = slide()
header(s, "04", "CATEGORY SPLIT", "分门类收益分成方案（建议）")
footer(s, idx)
rows = [
    ["门类", "收益内容", "甲方同普会", "乙方服务中心", "贡献逻辑"],
    ["载体运营补贴", "10–20万/年到平台", "40%", "60%", "乙方持资质运营；甲方协同认定与KPI"],
    ["孵化基地运营", "约10万/年", "35%", "65%", "乙方主体评估申报"],
    ["活动经费到账", "按主体到账部分", "50%", "50%", "甲方活动运营≈乙方申报资质"],
    ["奖项·申报服务费", "2–8万/家", "30%", "70%", "乙方执行申报；甲方导流获客"],
    ["奖项·成功分成", "到账额5%–15%", "40%", "60%", "共同促成兑现；乙方合规过审"],
    ["创新券核销分成", "服务机构核销", "30%", "70%", "乙方为入驻服务机构"],
    ["房租/专项成功费", "到账5%–15%", "40%", "60%", "甲方商务转化 + 乙方申报"],
]
yy = Inches(1.42)
cols = [0.14, 0.18, 0.12, 0.14, 0.42]
for ri, row in enumerate(rows):
    xx = ML
    h = Inches(0.58) if ri == 0 else Inches(0.62)
    for ci, val in enumerate(row):
        cw = Emu(int(CW * cols[ci]))
        bgc = RGBColor(0x2A, 0x1E, 0x55) if ri == 0 else (RGBColor(0x24, 0x1A, 0x4A) if ri % 2 else RGBColor(0x1C, 0x14, 0x3A))
        rect(s, xx, yy, cw, h, fill=bgc, line=LINE, lw=0.6)
        col = GOLD if ri == 0 else (GREEN if ci in (2, 3) and ri else INK)
        text(s, Emu(xx + Inches(0.08)), Emu(yy + Inches(0.1)), Emu(cw - Inches(0.12)), Emu(h - Inches(0.14)),
             str(val), size=11 if ri else 11.5, color=col, bold=(ri == 0 or ci in (0, 2, 3)),
             anchor=MSO_ANCHOR.MIDDLE)
        xx = Emu(xx + cw)
    yy = Emu(yy + h)
text(s, ML, Inches(6.5), CW, Inches(0.3),
     "加权后综合分成约：甲方 36%–40%　｜　乙方 60%–64%。可作为默认结算带，单项可±5–10 个点协商。",
     size=12, color=GOLD, bold=True)

# ========== 09 贡献度模型 ==========
idx = nxt()
s = slide()
header(s, "05", "CONTRIBUTION MODEL", "贡献度模型 · 如何定比例")
footer(s, idx)
# weights
dims = [
    ("资质与牌照", "25%", "乙方为主", "服务中心挂牌、载体认定、创新券机构"),
    ("申报执行", "25%", "乙方为主", "材料、过审、对接处室、合规风控"),
    ("获客与转化", "20%", "甲方为主", "同浦汇招商、活动引流、企业转化"),
    ("资源与背书", "15%", "双方共担", "协会/高校/政府关系与品牌"),
    ("资金与垫付", "10%", "按实约定", "活动垫资、第三方成本垫付"),
    ("日常运营", "5%", "双方共担", "对账、客服、档案与复盘"),
]
for i, (t, w, who, desc) in enumerate(dims):
    y = Inches(1.45) + Inches(i * 0.8)
    rect(s, ML, y, CW, Inches(0.72),
         grad=[(0, "2A1E55", None), (100, "1E1542", None)], line=LINE, lw=0.8, radius=True)
    text(s, Emu(ML + Inches(0.25)), Emu(y + Inches(0.18)), Inches(2.4), Inches(0.4),
         t, size=14, color=INK, bold=True)
    text(s, Emu(ML + Inches(2.8)), Emu(y + Inches(0.18)), Inches(1.2), Inches(0.4),
         [(w, GOLD, True)], size=18)
    text(s, Emu(ML + Inches(4.2)), Emu(y + Inches(0.2)), Inches(2.2), Inches(0.4),
         who, size=13, color=ACC if "乙" in who else (GREEN if "甲" in who else MUT), bold=True)
    text(s, Emu(ML + Inches(6.5)), Emu(y + Inches(0.2)), Inches(5.5), Inches(0.4),
         desc, size=12.5, color=MUT)

# ========== 10 测算结果 ==========
idx = nxt()
s = slide()
header(s, "05", "CALCULATION", "按贡献比例测算 · 成熟年分配")
footer(s, idx)

# Use mid case 220万 as service center realized revenue
rows = [
    ["门类", "成熟年服务方收益", "甲方分成", "甲方金额", "乙方分成", "乙方金额"],
    ["载体运营补贴", "20万", "40%", "8.0万", "60%", "12.0万"],
    ["孵化基地运营", "10万", "35%", "3.5万", "65%", "6.5万"],
    ["活动经费（中位）", "45万", "50%", "22.5万", "50%", "22.5万"],
    ["奖项·服务费", "55万", "30%", "16.5万", "70%", "38.5万"],
    ["奖项·成功分成", "40万", "40%", "16.0万", "60%", "24.0万"],
    ["其他政策（创新券等）", "30万", "35%", "10.5万", "65%", "19.5万"],
    ["合计（成熟年）", "200万", "约38.5%", "77.0万", "约61.5%", "123.0万"],
]
yy = Inches(1.42)
cols = [0.20, 0.16, 0.12, 0.14, 0.12, 0.14]
# normalize - 6 cols sum to 0.88? fix to 1.0
cols = [0.22, 0.16, 0.12, 0.14, 0.14, 0.22]
for ri, row in enumerate(rows):
    xx = ML
    h = Inches(0.52) if ri == 0 else Inches(0.55)
    last = ri == len(rows) - 1
    for ci, val in enumerate(row):
        cw = Emu(int(CW * cols[ci]))
        if ri == 0:
            bgc = RGBColor(0x2A, 0x1E, 0x55)
        elif last:
            bgc = RGBColor(0x3A, 0x2A, 0x18)
        else:
            bgc = RGBColor(0x24, 0x1A, 0x4A) if ri % 2 else RGBColor(0x1C, 0x14, 0x3A)
        rect(s, xx, yy, cw, h, fill=bgc, line=LINE, lw=0.6)
        col = GOLD if ri == 0 or last else INK
        text(s, Emu(xx + Inches(0.08)), Emu(yy + Inches(0.1)), Emu(cw - Inches(0.12)), Emu(h - Inches(0.12)),
             str(val), size=11.5 if last or ri == 0 else 11, color=col, bold=(ri == 0 or last or ci == 0),
             anchor=MSO_ANCHOR.MIDDLE)
        xx = Emu(xx + cw)
    yy = Emu(yy + h)
text(s, ML, Inches(6.35), CW, Inches(0.45),
     "解读：成熟年服务方到账约200万时，甲方同普会约77万、乙方服务中心约123万。若活动/奖项放量至275万，按同比例放大即可。",
     size=12, color=MUT)

# ========== 11 三年路径 ==========
idx = nxt()
s = slide()
header(s, "05", "3-YEAR PATH", "三年收益路径与累计分配（示意）")
footer(s, idx)
rows = [
    ["年度", "服务方到账（中位）", "甲方分成约", "乙方分成约", "阶段重点"],
    ["第1年 2027", "90–120万", "35–45万", "55–75万", "认定挂牌 · 样板3–5家 · 活动起步"],
    ["第2年 2028", "150–200万", "55–75万", "95–125万", "高企/专精特新放量 · 活动抬升"],
    ["第3年 2029", "200–275万", "75–105万", "125–170万", "双资质稳定 · 创新券核销常态"],
    ["三年累计", "440–595万", "165–225万", "275–370万", "综合分成约 甲38% / 乙62%"],
]
yy = Inches(1.55)
cols = [0.16, 0.18, 0.14, 0.14, 0.38]
for ri, row in enumerate(rows):
    xx = ML
    h = Inches(0.72)
    last = ri == len(rows) - 1
    for ci, val in enumerate(row):
        cw = Emu(int(CW * cols[ci]))
        if ri == 0:
            bgc = RGBColor(0x2A, 0x1E, 0x55)
        elif last:
            bgc = RGBColor(0x3A, 0x2A, 0x18)
        else:
            bgc = RGBColor(0x24, 0x1A, 0x4A) if ri % 2 else RGBColor(0x1C, 0x14, 0x3A)
        rect(s, xx, yy, cw, h, fill=bgc, line=LINE, lw=0.7)
        col = GOLD if ri == 0 or last else INK
        text(s, Emu(xx + Inches(0.1)), Emu(yy + Inches(0.18)), Emu(cw - Inches(0.16)), Emu(h - Inches(0.25)),
             str(val), size=13 if last or ri == 0 else 12.5, color=col, bold=(ri == 0 or last or ci == 0),
             anchor=MSO_ANCHOR.MIDDLE)
        xx = Emu(xx + cw)
    yy = Emu(yy + h)
text(s, ML, Inches(5.7), CW, Inches(0.9),
     [("结算机制　", GOLD, True),
      ("按自然季度对账；先扣除经双方确认的直接成本（第三方、物料、垫付），再按门类比例分配；"
       "一方对外收款的，到账后5个工作日内通知并划付。重大专项（如服务业引导资金项目制）单独立账、另议比例。", MUT, False)],
     size=13, space=1.2)

# ========== 12 结论 ==========
idx = nxt()
s = slide()
header(s, "06", "NEXT STEPS", "结论与签约建议")
footer(s, idx)
card(s, ML, Inches(1.5), Inches(5.7), Inches(5.1),
     title="结论（可直接对齐甲方）",
     items=[
         "服务中心成熟年中位收益约165–275万/年",
         "建议综合分成：甲方同普会约38% · 乙方约62%",
         "门类差异化：执行类乙方高、获客/活动类更均分",
         "奖项对企业侧是补贴，对服务方是服务费+分成",
         "数字为测算，写入协议时保留±协商带",
     ], accent=GOLD, bsize=13.5)
card(s, Inches(6.9), Inches(1.5), Inches(5.55), Inches(5.1),
     title="建议落地动作",
     items=[
         "① 签署《收益分成确认书》锁定门类比例",
         "② 完成服务中心挂牌 + 载体认定路径",
         "③ 入驻创新券平台，打通核销分成",
         "④ 锁定首批3–5家样板企业申报流水",
         "⑤ 季度联席对账，年度复盘调整比例",
         "⑥ 活动申报主体提前书面确认",
     ], accent=GREEN, bsize=13.5)

# fix footers total
total = len(_FOOT)
for run in _FOOT:
    run.text = " / %02d" % total

# save
os.makedirs(ART, exist_ok=True)
out_cn = os.path.join(HERE, "创智汇科创服务中心收益分成方案.pptx")
out_en = os.path.join(HERE, "chuangzhihui-service-center-revenue-share.pptx")
prs.save(out_cn)
prs.save(out_en)
import shutil
shutil.copy2(out_cn, os.path.join(ART, "创智汇科创服务中心收益分成方案.pptx"))
shutil.copy2(out_en, os.path.join(ART, "chuangzhihui-service-center-revenue-share.pptx"))
print("slides:", len(prs.slides))
print("Wrote", out_cn)
