# -*- coding: utf-8 -*-
"""生成「创智汇 6600㎡ AI+IP 产业创新中心合作方案」的 PPT 与 Excel 交付物。

用法: python3 build_ppt_excel.py
输出: 创智汇6600平合作方案.pptx / 创智汇6600平合作方案.xlsx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from openpyxl.worksheet.properties import PageSetupProperties

HERE = os.path.dirname(os.path.abspath(__file__))
FONT = "Microsoft YaHei"

# ---- palette ----
BG   = RGBColor(0x0B, 0x10, 0x20)
BG2  = RGBColor(0x16, 0x21, 0x3F)
INK  = RGBColor(0xEA, 0xF0, 0xFF)
MUT  = RGBColor(0x9F, 0xB0, 0xD6)
ACC  = RGBColor(0x4F, 0x8C, 0xFF)
ACC2 = RGBColor(0x7C, 0x5C, 0xFF)
GOLD = RGBColor(0xFF, 0xCE, 0x6B)
GREEN= RGBColor(0x3A, 0xD2, 0x9F)
LINE = RGBColor(0x2A, 0x37, 0x5A)

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s


def box(s, x, y, w, h, fill=None, line=None, radius=False):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    b = s.shapes.add_shape(shp_type, x, y, w, h)
    if fill is None:
        b.fill.background()
    else:
        b.fill.solid(); b.fill.fore_color.rgb = fill
    if line is None:
        b.line.fill.background()
    else:
        b.line.color.rgb = line; b.line.width = Pt(1)
    b.shadow.inherit = False
    return b


def text(s, x, y, w, h, runs, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, space=1.0):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    if isinstance(runs, str):
        runs = [(runs, color, bold)]
    p = tf.paragraphs[0]; p.alignment = align; p.line_spacing = space
    for t, c, b in runs:
        r = p.add_run(); r.text = t
        r.font.size = Pt(size); r.font.bold = b
        r.font.color.rgb = c; r.font.name = FONT
    return tb


def bullets(s, x, y, w, h, items, size=14, color=MUT, gap=6):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.12
        r = p.add_run(); r.text = "•  " + it
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = FONT
    return tb


def rule(s, x, y, w=Inches(1.0)):
    b = box(s, x, y, w, Pt(4), fill=ACC, radius=True)
    return b


def header(s, eyebrow, title, color=INK):
    text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.4),
         eyebrow, size=13, color=MUT, bold=True)
    text(s, Inches(0.7), Inches(0.85), Inches(12), Inches(0.7),
         title, size=30, color=color, bold=True)
    rule(s, Inches(0.72), Inches(1.62), Inches(1.4))


def table(s, x, y, w, rows, col_w, head_fill=BG2, sizes=None, head_color=GOLD):
    """rows[0] = header. col_w = list of fractional widths summing to 1."""
    n = len(rows); ncol = len(rows[0])
    row_h = Inches(0.52)
    sizes = sizes or [13]*ncol
    cy = y
    for ri, row in enumerate(rows):
        cx = x
        rh = row_h if ri == 0 else Inches(0.5)
        if ri % 2 == 1 and ri > 0:
            box(s, x, cy, w, rh, fill=RGBColor(0x12, 0x1B, 0x34))
        for ci, cell in enumerate(row):
            cw = Emu(int(w * col_w[ci]))
            col = head_color if ri == 0 else INK
            bold = ri == 0
            tb = s.shapes.add_textbox(cx, cy, cw, rh)
            tf = tb.text_frame; tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Pt(8); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
            p = tf.paragraphs[0]; p.line_spacing = 1.0
            r = p.add_run(); r.text = str(cell)
            r.font.size = Pt(sizes[ci] if ri > 0 else 12.5)
            r.font.bold = bold
            r.font.color.rgb = col if ci == 0 or ri == 0 else MUT if ci > 0 else INK
            r.font.name = FONT
            cx = Emu(cx + cw)
        box(s, x, Emu(cy + rh - Pt(1)), w, Pt(1), fill=LINE)
        cy = Emu(cy + rh)
    return cy


def card(s, x, y, w, h, title, body=None, items=None, accent=ACC, title_color=INK):
    box(s, x, y, w, h, fill=BG2, line=LINE, radius=True)
    box(s, Emu(x+Inches(0.18)), Emu(y+Inches(0.2)), Inches(0.07), Inches(0.32), fill=accent)
    text(s, Emu(x+Inches(0.36)), Emu(y+Inches(0.16)), Emu(w-Inches(0.5)), Inches(0.4),
         title, size=15, color=title_color, bold=True)
    if body:
        text(s, Emu(x+Inches(0.28)), Emu(y+Inches(0.62)), Emu(w-Inches(0.5)), Emu(h-Inches(0.7)),
             body, size=12.5, color=MUT, space=1.12)
    if items:
        bullets(s, Emu(x+Inches(0.28)), Emu(y+Inches(0.6)), Emu(w-Inches(0.5)), Emu(h-Inches(0.7)),
                items, size=12.5)


def footer(s, idx):
    text(s, Inches(0.7), Inches(7.05), Inches(9), Inches(0.3),
         "创智汇 6600㎡ · AI+IP 产业创新中心招商运营合作 · 谈判集中汇报", size=10, color=MUT)
    text(s, Inches(11.6), Inches(7.05), Inches(1.4), Inches(0.3),
         str(idx), size=10, color=MUT, align=PP_ALIGN.RIGHT)


# ============ SLIDES ============
# 1 cover
s = slide()
box(s, 0, 0, SW, SH, fill=BG)
box(s, Inches(0), Inches(0), Inches(0.18), SH, fill=ACC)
text(s, Inches(0.9), Inches(1.3), Inches(11.5), Inches(0.4),
     "CHUANGZHIHUI · WUJIAOCHANG URBAN RENEWAL", size=14, color=MUT, bold=True)
text(s, Inches(0.9), Inches(1.9), Inches(11.6), Inches(1.6),
     [("创智汇 ", INK, True), ("6600㎡", GOLD, True), (" · AI + IP 产业创新中心", INK, True),
      ("\n招商运营合作 — 谈判集中汇报", INK, True)], size=38, bold=True, space=1.1)
rule(s, Inches(0.92), Inches(3.7), Inches(1.6))
text(s, Inches(0.9), Inches(3.95), Inches(11.6), Inches(1.0),
     [("五角场片区城市更新 · 创智汇（一期）OPC + AI + IP 创新中心\n", MUT, False),
      ("3 楼 孵化器 + 办公 ≈ ", MUT, False), ("2850㎡", GOLD, True),
      ("    ·    5 楼 展厅 + 贸易 ≈ ", MUT, False), ("3670㎡", GOLD, True)],
     size=17, space=1.4)
for i, t in enumerate(["上海·杨浦·五角场·创智天地", "环同济经济圈·超级链接器",
                       "IP 为内容·AI 为工具·空间为载体", "参照「森马×元谷」打法编制"]):
    bx = Inches(0.9 + i*3.0)
    box(s, bx, Inches(5.3), Inches(2.85), Inches(0.5), fill=BG2, line=LINE, radius=True)
    text(s, bx, Inches(5.3), Inches(2.85), Inches(0.5), t, size=11.5, color=INK,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(0.9), Inches(6.4), Inches(11), Inches(0.4),
     "v1.0 谈判稿 · 本汇报仅用于商务谈判沟通", size=12, color=MUT)

# 2 agenda
s = slide(); header(s, "AGENDA", "汇报议程 · 七大板块")
ag = [("00", "关键信息", "材料浓缩为事实底盘"), ("01", "项目与共识", "定位/双重身份/谈判要点"),
      ("02", "我们带来什么", "六大资源资产化"), ("03", "招商方案", "IP+AI 双轨/楼层/漏斗"),
      ("04", "节点与排期", "启动-导入-成型-提升"), ("05", "品牌与活动", "挂牌/沙龙/漫展/会客厅"),
      ("06", "条款与收费", "两模式+五类收费+测算"), ("07", "投决与下一步", "三问决策+30/60/90 天")]
cw, ch = Inches(2.95), Inches(1.5)
for i, (n, t, d) in enumerate(ag):
    x = Inches(0.7 + (i % 4)*3.1); y = Inches(2.0 + (i // 4)*1.75)
    box(s, x, y, cw, ch, fill=BG2, line=LINE, radius=True)
    box(s, Emu(x+Inches(0.22)), Emu(y+Inches(0.22)), Inches(0.55), Inches(0.5), fill=ACC2, radius=True)
    text(s, Emu(x+Inches(0.22)), Emu(y+Inches(0.22)), Inches(0.55), Inches(0.5), n,
         size=15, color=BG, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Emu(x+Inches(0.9)), Emu(y+Inches(0.26)), Inches(2.0), Inches(0.5), t, size=15, color=INK, bold=True)
    text(s, Emu(x+Inches(0.24)), Emu(y+Inches(0.85)), Inches(2.5), Inches(0.6), d, size=12, color=MUT)
footer(s, 2)

# 3 fact sheet
s = slide(); header(s, "00 · FACT SHEET", "关键信息提取 · 项目基本盘")
y0 = Inches(1.9)
box(s, Inches(0.7), y0, Inches(6.0), Inches(3.0), fill=BG2, line=LINE, radius=True)
table(s, Inches(0.85), Emu(y0+Inches(0.18)), Inches(5.7), [
    ["项目", "内容"], ["位置", "上海杨浦五角场·创智汇"], ["合作面积", "约 6600㎡"],
    ["3 楼 ≈2850㎡", "孵化器+办公（AI/OPC 主轴）"], ["5 楼 ≈3670㎡", "展厅+贸易（IP 内容主轴）"],
    ["属性", "产业空间+科创合作+招商运营"]], [0.32, 0.68], sizes=[12.5, 12.5])
card(s, Inches(6.95), y0, Inches(5.65), Inches(3.0), "区位与资源优势", items=[
    "环同济：复旦/同济/财大/上理工高校群",
    "TOD + 五角场商圈成熟，通勤配套俱佳",
    "杨浦三区联动 + 五角场五区联动创新高地",
    "离我方极近(20–30min)→支撑「轻驻场·高频次」运营"], accent=GOLD)
yy = Inches(5.15)
card(s, Inches(0.7), yy, Inches(3.85), Inches(1.6), "定位主线",
     body="超级链接器：IP 为内容、AI 为工具、空间为载体；产业筋骨·文化灵魂·商业血脉。")
card(s, Inches(4.7), yy, Inches(3.85), Inches(1.6), "三大集群",
     body="动漫 IP · 科技应用 · 交互设计", accent=GREEN)
card(s, Inches(8.7), yy, Inches(3.9), Inches(1.6), "可链接资源",
     body="北大上海校友会、同济设计创新院、科企联、IP/玩具/广告协会、混知等 IP、聚成智能、中建四局。", accent=ACC2)
footer(s, 3)

# 4 section 01 + consensus
s = slide(); header(s, "01 · 项目与共识", "双重身份 + 三件谈判要点")
card(s, Inches(0.7), Inches(1.9), Inches(5.95), Inches(1.5), "身份一 · 产业孵化核心",
     body="五角场城市更新的孵化核心与先行示范区——空间改造+产业升级+商业激活的完整落地路径。")
card(s, Inches(6.85), Inches(1.9), Inches(5.75), Inches(1.5), "身份二 · 超级链接器",
     body="环同济经济圈资源链接平台：IP 内容 × AI 工具 × 空间载体，闭环生长。", accent=ACC2)
text(s, Inches(0.7), Inches(3.6), Inches(8), Inches(0.4), "本次谈判要点", size=18, color=INK, bold=True)
pts = [("1", "合作模式", "轻运营(不对赌) 还是 对赌运营(对赌招商)？"),
       ("2", "商业条款", "月费+招商佣金+挂牌+活动+媒体 的组合与金额"),
       ("3", "启动时间", "与装修/设计交付同步，尽快锁定并启动招商")]
for i, (n, t, d) in enumerate(pts):
    x = Inches(0.7 + i*4.0)
    box(s, x, Inches(4.1), Inches(3.8), Inches(1.4), fill=BG2, line=LINE, radius=True)
    box(s, Emu(x+Inches(0.22)), Emu(Inches(4.1)+Inches(0.22)), Inches(0.5), Inches(0.5), fill=ACC, radius=True)
    text(s, Emu(x+Inches(0.22)), Emu(Inches(4.1)+Inches(0.22)), Inches(0.5), Inches(0.5), n, size=15, color=BG, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Emu(x+Inches(0.85)), Emu(Inches(4.1)+Inches(0.26)), Inches(2.7), Inches(0.4), t, size=15, color=INK, bold=True)
    text(s, Emu(x+Inches(0.24)), Emu(Inches(4.1)+Inches(0.82)), Inches(3.4), Inches(0.5), d, size=11.5, color=MUT)
box(s, Inches(0.7), Inches(5.75), Inches(11.9), Inches(1.05), fill=RGBColor(0x14, 0x20, 0x40), line=GOLD, radius=True)
text(s, Inches(0.95), Inches(5.85), Inches(11.5), Inches(0.9),
     [("核心共识：", GOLD, True), ("把 6600㎡ 升级为「产业空间合作运营项目」——我方输出 政策+招商+产业资源+活动+媒体 五项能力，按「基础服务费+招商佣金+增值分成」取酬。", INK, False)],
     size=14, anchor=MSO_ANCHOR.MIDDLE, space=1.2)
footer(s, 4)

# 5 six resources
s = slide(); header(s, "02 · 六大资源 → 招商引擎", "资源即招商")
res = [("① 政策抓手", "杨浦 AI/大数据补贴存量续享 + 算力补贴最高 50% + 腾讯云 85 折/免 2 月 →「创智汇专有政策包」", GOLD),
       ("② 高校/协会背书", "北大上海校友会、同济设计创新院、复旦、科企联、IP/玩具/广告协会 → 挂牌即招商", ACC),
       ("③ IP 内容资源", "混知、观山礼玉等 IP 品牌方 + 汕头玩具/扬州毛绒/东莞潮玩/文创集群 → 5 楼自带内容", ACC2),
       ("④ AI/OPC 社群", "自有 OPC(AI 开放社群)品牌、黑客松、AI 项目库、聚成智能 → 3 楼自带客流", GREEN),
       ("⑤ 活动/会客厅", "漫展(门票+赞助)、北欧会客厅(出海撮合)、沙龙、峰会借势 → 活动即招商", ACC),
       ("⑥ 资本/产业方", "江西金控基金、景德镇陶瓷版权交易中心、中建四局工程资源、奇瑞/华为外延 → 链主带动", GOLD)]
for i, (t, d, c) in enumerate(res):
    x = Inches(0.7 + (i % 3)*4.0); y = Inches(2.0 + (i // 3)*1.95)
    card(s, x, y, Inches(3.8), Inches(1.75), t, body=d, accent=c)
text(s, Inches(0.7), Inches(6.05), Inches(12), Inches(0.5),
     [("普通园区「发广告·等上门」；六维齐发 → 客户 ", MUT, False), ("「被送进来」而非「被拉进来」。", GOLD, True)], size=14)
footer(s, 5)

# 6 floor split
s = slide(); header(s, "03 · 双轨与楼层定位", "3 楼孵化办公(AI 主轴) + 5 楼展厅贸易(IP 主轴)")
box(s, Inches(0.7), Inches(1.9), Inches(5.95), Inches(4.6), fill=BG2, line=LINE, radius=True)
text(s, Inches(0.95), Inches(2.05), Inches(5.6), Inches(0.4), "3 楼 · 约 2850㎡ · 孵化+办公", size=16, color=ACC, bold=True)
table(s, Inches(0.9), Inches(2.55), Inches(5.6), [
    ["产品", "面积", "客户"], ["标准小单元", "80–150㎡×8–12", "AI 初创/设计/小微"],
    ["成长型单元", "200–350㎡×3–4", "AI 应用/研发服务"], ["OPC 联合办公", "工位大区", "AI 项目/黑客松/孵化"],
    ["直播/AI 展示", "共享 1–2 间", "按次/时段"]], [0.33, 0.34, 0.33], sizes=[12, 11.5, 11.5])
box(s, Inches(6.85), Inches(1.9), Inches(5.75), Inches(4.6), fill=BG2, line=LINE, radius=True)
text(s, Inches(7.1), Inches(2.05), Inches(5.4), Inches(0.4), "5 楼 · 约 3670㎡ · 展厅+贸易", size=16, color=GOLD, bold=True)
table(s, Inches(7.05), Inches(2.55), Inches(5.4), [
    ["产品", "面积", "内容"], ["综合集群展厅", "700–900㎡", "50+ IP 轮展"],
    ["产业集群展位", "270–700㎡×4–5", "玩具/毛绒/潮玩/文创"], ["IP 展销/跨境", "按摊位", "IP 零售/跨境展销"],
    ["培训沙龙/仓储", "120–650㎡", "OPC 培训/活动"]], [0.33, 0.34, 0.33], sizes=[12, 11.5, 11.5])
footer(s, 6)

# 7 funnel + targets
s = slide(); header(s, "03 · 招商漏斗与务实目标", "四级漏斗 + 务实去化")
fn = [("L1", "牌照锚定", "校友会/科企联/IP/玩具协会挂牌前置，客户送进来"),
      ("L2", "政策招商", "AI/大数据续享+算力补贴+腾讯云作「入驻礼包」"),
      ("L3", "社群带流", "OPC 社群、黑客松、小红书投流、漫展、沙龙持续导流"),
      ("L4", "资源转化", "高校成果转化、IP 品牌方、链主配套、外贸玩具厂转化")]
for i, (n, t, d) in enumerate(fn):
    x = Inches(0.7 + i*3.05)
    box(s, x, Inches(2.0), Inches(2.85), Inches(1.9), fill=BG2, line=LINE, radius=True)
    text(s, Emu(x+Inches(0.25)), Inches(2.2), Inches(1.5), Inches(0.5), n, size=20, color=ACC, bold=True)
    text(s, Emu(x+Inches(0.25)), Inches(2.75), Inches(2.4), Inches(0.4), t, size=15, color=INK, bold=True)
    text(s, Emu(x+Inches(0.25)), Inches(3.2), Inches(2.45), Inches(0.7), d, size=11.5, color=MUT, space=1.1)
kpi = [("2000–3000㎡", "一年去化达标线(不设过高保底)", GOLD),
       ("3300–4300㎡", "对赌模式建议锁定对赌面积", ACC),
       ("5 类客群", "AI/潮玩玩具文创 IP/科技中小/高校转化/生产性服务", GREEN)]
for i, (v, l, c) in enumerate(kpi):
    x = Inches(0.7 + i*4.0)
    box(s, x, Inches(4.4), Inches(3.8), Inches(1.7), fill=BG2, line=LINE, radius=True)
    text(s, Emu(x+Inches(0.3)), Inches(4.65), Inches(3.4), Inches(0.6), v, size=26, color=c, bold=True)
    text(s, Emu(x+Inches(0.3)), Inches(5.4), Inches(3.4), Inches(0.6), l, size=12, color=MUT, space=1.1)
footer(s, 7)

# 8 milestones
s = slide(); header(s, "04 · 四阶段排期", "从定位包装到稳定运营")
table(s, Inches(0.7), Inches(2.0), Inches(11.9), [
    ["阶段", "时间窗口", "目标", "重点动作"],
    ["启动期", "0–3 月", "定位+政策+渠道", "招商手册、收费标准、政策汇编、企业库、首批挂牌、1–2 家样板"],
    ["导入期", "3–6 月", "首批入驻 30–50%", "集中招引 AI/IP/潮玩、首批优惠、推介路演、政策诊断"],
    ["成型期", "6–12 月", "提升出租率与质量", "重点客户补位、专精特新/高企培育、企业服务收费、首场漫展"],
    ["提升期", "12 月+", "稳定收入与品牌", "区级示范点、载体资质、活动 IP 化、出海撮合、模式复制"]],
    [0.12, 0.14, 0.2, 0.54], sizes=[13, 12.5, 12.5, 12])
footer(s, 8)

# 9 brand & events
s = slide(); header(s, "05 · 品牌与活动", "活动即招商，品牌即势能")
card(s, Inches(0.7), Inches(2.0), Inches(3.85), Inches(2.2), "挂牌 · 牌照即招商",
     body="北大上海校友会(优先)、科企联/科技服务中心、IP 协会、中国玩具协会。\n\n总包约 20–50 万", accent=GOLD)
card(s, Inches(4.7), Inches(2.0), Inches(3.85), Inches(2.2), "活动 · 全年节奏",
     body="20–24 场沙龙/路演/培训/发布会。\n\n打包约 30 万(+赞助/门票)", accent=ACC)
card(s, Inches(8.7), Inches(2.0), Inches(3.9), Inches(2.2), "漫展 · 自造流量",
     body="千人级，门票+赞助；10–11 月首场，对接 BW 广告商与二次元资源。", accent=ACC2)
box(s, Inches(0.7), Inches(4.5), Inches(11.9), Inches(1.7), fill=BG2, line=LINE, radius=True)
text(s, Inches(0.95), Inches(4.65), Inches(11.5), Inches(0.4), "北欧会客厅 · 出海撮合", size=15, color=GREEN, bold=True)
text(s, Inches(0.95), Inches(5.15), Inches(11.4), Inches(0.95),
     [("复用品牌势能做海外 IP/技术入华与国内 IP 出海撮合台。取酬二选一：", MUT, False),
      ("收租金则不分出海成交；不收租金、企业自费布展则在撮合成交中取分成。", INK, True)], size=13.5, space=1.25)
footer(s, 9)

# 10 two modes
s = slide(); header(s, "06 · 两种合作模式", "先定大方向：轻运营 vs 对赌运营")
box(s, Inches(0.7), Inches(2.0), Inches(5.95), Inches(3.4), fill=BG2, line=LINE, radius=True)
text(s, Inches(0.95), Inches(2.15), Inches(5.6), Inches(0.4), "模式 A · 轻运营(不对赌招商)", size=16, color=INK, bold=True)
bullets(s, Inches(0.95), Inches(2.7), Inches(5.5), Inches(2.5),
        ["不驻场、不背去化指标", "收费：活动费+挂牌费+媒体策划费", "适用：对方只要内容/活动/品牌"], size=14, gap=10)
box(s, Inches(6.85), Inches(2.0), Inches(5.75), Inches(3.4), fill=RGBColor(0x1A, 0x22, 0x40), line=GOLD, radius=True)
text(s, Inches(7.1), Inches(2.15), Inches(5.4), Inches(0.4), "模式 B · 对赌运营(推荐)", size=16, color=GOLD, bold=True)
bullets(s, Inches(7.1), Inches(2.7), Inches(5.3), Inches(2.5),
        ["轻驻场(1 人/接待位，离得近、可挂牌办公)+背对赌去化", "收费：基础月费+招商佣金+活动+挂牌+媒体", "适用：对方要招商结果、愿付运营费"], size=14, gap=10)
text(s, Inches(0.7), Inches(5.7), Inches(12), Inches(0.6),
     [("谈判原则：", MUT, False), ("先按高位报，再让对方还价", GOLD, True), ("；收月费即对赌去化，不收月费则不背指标。", MUT, False)], size=14)
footer(s, 10)

# 11 fee table
s = slide(); header(s, "06 · 五类收费项(收费项建议)", "收费清单")
table(s, Inches(0.7), Inches(1.9), Inches(11.9), [
    ["类别", "对象", "方式", "建议金额 / 口径"],
    ["① 基础运营月费", "合作方", "按月", "市场 10–30 万；首报 12 万；对赌版 5–6 万(对赌 3300–4300㎡)"],
    ["② 招商佣金", "合作方/业主", "成交后", "市场 2–3 月年租金；我方抽 ≤1 月，只给首月不重复"],
    ["③ 挂牌费", "合作方", "一次性", "总包 20–50 万(按数量与影响力分档)"],
    ["④ 活动执行费", "合作方+外部", "按场/打包", "20–24 场/年约 30 万 + 赞助/门票分成"],
    ["⑤ 媒体流量费", "合作方", "按季/项目", "小红书投流、OPC 内容、直播间运营"],
    ["⑥ 企业增值服务", "入驻企业", "按项/成功", "工商注册、政策申报、知识产权、财税法务、融资、出海撮合"]],
    [0.2, 0.13, 0.13, 0.54], sizes=[12.5, 12, 12, 11.5])
footer(s, 11)

# 12 package + math
s = slide(); header(s, "06 · 打包报价 + 对方视角测算", "算得过账：花一份钱，换更快去化与更优企业")
box(s, Inches(0.7), Inches(2.0), Inches(5.95), Inches(4.2), fill=BG2, line=LINE, radius=True)
text(s, Inches(0.95), Inches(2.15), Inches(5.6), Inches(0.4), "打包报价示例(谈判锚)", size=15, color=INK, bold=True)
bullets(s, Inches(0.95), Inches(2.7), Inches(5.5), Inches(3.4), [
    "固定服务包：月费 5 万×12=60 万 + 活动 30 万 + 挂牌 20 万 = 约 110 万/年(可砍至~100 万)",
    "浮动：招商佣金按去化另计(我方抽 ≤1 月)；驻场/媒体/出海分成据实另计",
    "政策申报：普通 0.3–2 万/项；高企/专精特新 2–8 万/项；补贴类按到账 5%–15% 成功费"], size=13, gap=12)
box(s, Inches(6.85), Inches(2.0), Inches(5.75), Inches(4.2), fill=BG2, line=LINE, radius=True)
text(s, Inches(7.1), Inches(2.15), Inches(5.4), Inches(0.4), "对方视角测算", size=15, color=INK, bold=True)
table(s, Inches(7.05), Inches(2.65), Inches(5.4), [
    ["维度", "口径"], ["对方空间成本", "≈554 万/年(2.3 元/㎡/天×6600×365)"],
    ["我方服务包", "100–120 万/年 + 招商佣金(按去化)"], ["列支来源", "对方利润/中建四局运营费，非纯增量"],
    ["对方所得", "去化↑、企业质量↑、政策承接、品牌示范"]], [0.32, 0.68], sizes=[12.5, 12])
footer(s, 12)

# 13 decision
s = slide(); header(s, "07 · 投决建议与下一步", "三问 + 启动路线图")
q = [("Q1 合作模式？", "建议 模式 B 轻量版：离得近→轻驻场+对赌去化，盈利性最佳", GOLD),
     ("Q2 月费档位？", "首报 12 万/月(市场 10–30 万)，对赌版可落 5–6 万；打包年费锚 100–120 万", ACC),
     ("Q3 何时启动？", "与装修/设计交付同步，尽快签 MOU、先行政策梳理与挂牌筹备", ACC2)]
for i, (t, d, c) in enumerate(q):
    x = Inches(0.7 + i*4.0)
    card(s, x, Inches(1.95), Inches(3.8), Inches(1.55), t, body=d, accent=c)
road = [("30", "天", "签 MOU + 招商手册与收费方案 + 政策汇编 + 首批挂牌对接"),
        ("60", "天", "首场沙龙/路演 + OPC 社群与投流 + 首单招商 + 漫展筹备"),
        ("90", "天", "正式挂牌 + 样板企业入驻 + 去化向 30–50% 推进 + 漫展锁档")]
for i, (n, u, d) in enumerate(road):
    x = Inches(0.7 + i*4.0)
    box(s, x, Inches(3.7), Inches(3.8), Inches(1.55), fill=BG2, line=LINE, radius=True)
    text(s, Emu(x+Inches(0.25)), Inches(3.85), Inches(2), Inches(0.5), [(n, ACC, True), (" "+u, MUT, False)], size=20, bold=True)
    text(s, Emu(x+Inches(0.25)), Inches(4.4), Inches(3.4), Inches(0.8), d, size=12, color=MUT, space=1.15)
box(s, Inches(0.7), Inches(5.5), Inches(11.9), Inches(1.0), fill=RGBColor(0x14, 0x20, 0x40), line=GOLD, radius=True)
text(s, Inches(0.95), Inches(5.6), Inches(11.5), Inches(0.8),
     [("分步策略：", GOLD, True), ("先签 6600㎡ 单项目专项(单项目单核算)，跑通后再承接体系内其他物业与外延项目。", INK, False)],
     size=14, anchor=MSO_ANCHOR.MIDDLE, space=1.2)
footer(s, 13)

# 14 thanks
s = slide()
box(s, 0, 0, SW, SH, fill=BG)
box(s, Inches(0), Inches(0), Inches(0.18), SH, fill=ACC)
text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.4), "THANKS", size=15, color=MUT, bold=True)
text(s, Inches(0.9), Inches(2.9), Inches(11.5), Inches(1.8),
     [("让创智汇成为杨浦五角场\n", INK, True), ("AI + IP 产业与城市更新", ACC, True), ("的超级链接器与可复制样板", INK, True)],
     size=34, bold=True, space=1.2)
text(s, Inches(0.9), Inches(4.9), Inches(11.5), Inches(0.5),
     "90 天启动 · 政策借势 · 闭环招商 · 内容自造流量", size=18, color=MUT)
text(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.4),
     "谈判集中汇报 v1.0 · 期待达成合作", size=12, color=MUT)

ppt_path = os.path.join(HERE, "创智汇6600平合作方案.pptx")
prs.save(ppt_path)
print("PPT saved:", ppt_path, "slides:", len(prs.slides._sldIdLst))


# ============ EXCEL ============
wb = Workbook()
HEAD_FILL = PatternFill("solid", fgColor="16213F")
TITLE_FILL = PatternFill("solid", fgColor="0B1020")
ALT_FILL = PatternFill("solid", fgColor="F2F5FC")
thin = Side(style="thin", color="C9D4EC")
BORDER = Border(left=thin, right=thin, top=thin, bottom=thin)
WHITE = Font(name=FONT, color="FFFFFF", bold=True, size=11)
GOLDF = Font(name=FONT, color="B8860B", bold=True)
TITLEF = Font(name=FONT, color="FFFFFF", bold=True, size=14)
NORM = Font(name=FONT, size=10.5)
WRAP = Alignment(wrap_text=True, vertical="center")
CEN = Alignment(horizontal="center", vertical="center", wrap_text=True)


def sheet(name, title, headers, rows, widths, title_span=None):
    ws = wb.create_sheet(name)
    ncol = len(headers)
    span = title_span or ncol
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=span)
    c = ws.cell(1, 1, title); c.font = TITLEF; c.fill = TITLE_FILL
    c.alignment = Alignment(vertical="center", horizontal="left", indent=1)
    ws.row_dimensions[1].height = 32
    for j, h in enumerate(headers, 1):
        cc = ws.cell(2, j, h); cc.font = WHITE; cc.fill = HEAD_FILL
        cc.alignment = CEN; cc.border = BORDER
    ws.row_dimensions[2].height = 24
    for i, row in enumerate(rows, 3):
        for j, val in enumerate(row, 1):
            cc = ws.cell(i, j, val); cc.font = NORM; cc.alignment = WRAP; cc.border = BORDER
            if i % 2 == 1:
                cc.fill = ALT_FILL
        ws.row_dimensions[i].height = 30
    for j, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(j)].width = w
    ws.sheet_view.showGridLines = False
    ws.freeze_panes = "A3"
    # 打印/导出：横向 + 缩放至一页宽，避免列被拆到多页
    ws.page_setup.orientation = "landscape"
    ws.page_setup.fitToWidth = 1
    ws.page_setup.fitToHeight = 0
    ws.sheet_properties.pageSetUpPr = PageSetupProperties(fitToPage=True)
    ws.print_options.horizontalCentered = True
    ws.page_margins.left = ws.page_margins.right = 0.3
    ws.page_margins.top = ws.page_margins.bottom = 0.4
    return ws


wb.remove(wb.active)

sheet("01 关键信息", "00 · 关键信息提取 / 项目基本盘",
      ["项目", "内容"],
      [["项目位置", "上海市杨浦区五角场附近 · 创智汇(创智天地片区)"],
       ["合作面积", "约 6600㎡"],
       ["3 楼(约 2850㎡)", "孵化器 + 办公（AI / OPC 主轴）"],
       ["5 楼(约 3670㎡)", "展厅 + 贸易（IP 内容主轴）"],
       ["项目属性", "产业空间 + 科创合作空间 + 企业招商运营空间"],
       ["合作方向", "产业招商、空间运营、企业服务、政策服务、IP 内容与出海"],
       ["核心价值", "区位 + 杨浦科创资源 + 高校人才 + 产业政策 + 城市更新空间价值"],
       ["定位主线", "超级链接器：IP 为内容、AI 为工具、空间为载体"],
       ["三大集群", "动漫 IP、科技应用、交互设计"],
       ["可链接资源", "北大上海校友会、同济设计创新院、科企联、IP/玩具/广告协会、混知等 IP、聚成智能、中建四局"],
       ["关键变量", "离我方极近(20–30 分钟)→ 支撑「轻驻场、高频次」运营，是盈利关键"]],
      [22, 78])

sheet("02 楼层产品", "03 · 楼层产品切分（招商最小颗粒度）",
      ["楼层", "产品", "面积区间", "数量建议", "适合客户"],
      [["3 楼 孵化+办公\n≈2850㎡(AI/OPC 主轴)", "标准小单元", "80–150㎡", "8–12 个", "AI 初创、设计团队、科技型小微"],
       ["", "成长型单元", "200–350㎡", "3–4 个", "AI 应用、研发设计、生产性服务"],
       ["", "OPC 联合办公/工位", "工位大区", "1 个大区", "AI 项目、黑客松团队、孵化项目"],
       ["", "直播间 / AI 展示运营", "共享", "1–2 间", "按次/按时段共享"],
       ["5 楼 展厅+贸易\n≈3670㎡(IP 内容主轴)", "综合集群展厅", "700–900㎡", "1 个", "50+ IP 集中展示、轮展"],
       ["", "产业集群展位", "270–700㎡/区", "4–5 区", "汕头玩具/扬州毛绒/东莞潮玩/文创"],
       ["", "IP 展销 / 跨境贸易", "按摊位/面积", "若干", "IP 零售、跨境展销平台"],
       ["", "培训沙龙 / 仓储", "120–650㎡", "各 1", "OPC 社区培训、活动、样品仓储"]],
      [22, 20, 16, 12, 32])

sheet("03 招商策略", "03 · 招商漏斗 + 务实目标",
      ["层级/项", "名称", "说明"],
      [["L1", "牌照锚定", "校友会/科企联/IP/玩具协会挂牌前置，客户「送进来」"],
       ["L2", "政策招商", "AI/大数据补贴续享 + 算力补贴 + 腾讯云作「入驻礼包」"],
       ["L3", "社群/活动带流", "OPC 社群、黑客松、小红书投流、漫展、沙龙持续导流"],
       ["L4", "资源转化", "高校成果转化、IP 品牌方、链主配套、外贸玩具厂转化"],
       ["目标", "年去化达标线", "2000–3000㎡（不设过高保底）"],
       ["目标", "对赌面积", "对赌模式建议锁定 3300–4300㎡"],
       ["客群", "五类目标客户", "AI 项目 / 潮玩玩具文创 IP / 科技型中小 / 高校成果转化 / 生产性服务业"]],
      [12, 18, 70])

sheet("04 节点排期", "04 · 四阶段排期",
      ["阶段", "时间窗口", "目标", "重点动作"],
      [["启动期", "0–3 月", "定位+政策+渠道", "招商手册、收费标准、政策汇编、企业库、首批挂牌、1–2 家样板"],
       ["导入期", "3–6 月", "首批入驻 30–50%", "集中招引 AI/IP/潮玩、首批优惠、推介路演、政策诊断"],
       ["成型期", "6–12 月", "提升出租率与质量", "重点客户补位、专精特新/高企培育、企业服务收费、首场漫展"],
       ["提升期", "12 月+", "稳定收入与品牌", "区级示范点、载体资质、活动 IP 化、出海撮合、模式复制"]],
      [12, 14, 22, 56])

sheet("05 品牌活动", "05 · 品牌与活动",
      ["板块", "内容", "金额 / 口径"],
      [["挂牌", "北大上海校友会(优先)、科企联/科技服务中心、IP 协会、中国玩具协会", "总包约 20–50 万"],
       ["活动", "20–24 场沙龙/路演/培训/发布会", "打包约 30 万(+赞助/门票)"],
       ["漫展", "千人级，门票+赞助；10–11 月首场，对接 BW 广告商与二次元资源", "门票+赞助双收入"],
       ["北欧会客厅", "海外 IP/技术入华 + 国内 IP 出海撮合台", "收租金则不分出海；不收租金则撮合成交取分成(二选一)"]],
      [16, 56, 34])

# core fee sheet
ws = sheet("06 收费项建议(核心)", "06 · 五类收费项（收费项建议）",
      ["类别", "收费对象", "收费方式", "建议金额 / 口径"],
      [["① 基础运营月费", "合作方/项目公司", "按月", "市场 10–30 万/月；首报 12 万/月；对赌版 5–6 万/月(对赌 3300–4300㎡)；收月费即对赌去化"],
       ["② 招商佣金", "合作方/业主方", "成交后", "市场 2–3 个月年租金；我方抽 ≤1 个月，只给首月、不重复；返投/重点客户可另议"],
       ["③ 挂牌费", "合作方", "一次性/按项", "总包 20–50 万(北大校友会/科企联/IP 协会等，按数量与影响力分档)"],
       ["④ 活动执行费", "合作方+外部", "按场/打包", "20–24 场/年打包约 30 万；漫展门票+赞助、净利分成"],
       ["⑤ 媒体流量费", "合作方", "按季度/项目", "小红书投流、OPC 社群内容、短视频/直播间运营"],
       ["⑥ 企业增值服务", "入驻企业", "按项/成功收费", "工商注册、政策申报(科小备案/高企/专精特新/研发归集/技改)、知识产权、财税法务、融资、出海撮合"]],
      [18, 18, 14, 50])

sheet("07 报价测算", "06 · 打包报价示例 + 对方视角测算",
      ["项目", "口径 / 金额", "说明"],
      [["固定服务包(年)", "约 110 万/年(可砍至~100 万)", "月费 5 万×12=60 万 + 活动 30 万 + 挂牌 20 万"],
       ["招商佣金(浮动)", "我方抽 ≤1 个月年租金", "按实际去化另计，只给首月不重复"],
       ["增值服务(浮动)", "据实另计", "政策申报、出海撮合分成、驻场运营费、媒体投流"],
       ["政策申报参考", "0.3–2 万/项；高企专精特新 2–8 万/项；补贴类到账额 5%–15%", "成功费按实际兑现计"],
       ["——对方测算——", "", ""],
       ["对方空间成本", "≈554 万/年", "2.3 元/㎡/天 × 6600㎡ × 365，一天约 1.5 万+"],
       ["我方服务包", "100–120 万/年 + 招商佣金", "由对方利润/中建四局运营费列支，非纯增量"],
       ["对方收益", "去化率↑、企业质量↑、政策承接、租金溢价、品牌示范", "比照元谷「花一份钱赚四份钱」逻辑，账面可算正"]],
      [22, 40, 38])

sheet("08 政策结合", "06 · 当地政策内容结合（以政府最新文件为准）",
      ["政策方向", "结合方式 / 抓手", "在本项目的用法"],
      [["杨浦 AI/大数据补贴", "存量延续：2025.2 到期，但已申报企业 3 年内可续享", "作为招商引子，引导企业向我方靠拢"],
       ["算力补贴", "凭算力发票向科委/经信申请，最高补贴 50%", "打造「创智汇专有政策包」"],
       ["腾讯云 / 算力合作", "新注册 AI 公司免 2 个月 / 算力 85 折", "签约共建，叠加专有政策"],
       ["科技型中小企业", "辅导备案，为高企/研发加计扣除打基础", "基础政策诊断(免费引流)"],
       ["高新技术企业", "研发投入/知识产权/科技人员比例辅导", "增值收费项"],
       ["专精特新", "创新型中小→市级专精特新→国家级小巨人", "固定费+成功费"],
       ["人才政策", "落户、人才公寓、创业人才、毕业生招聘", "企业服务包"],
       ["研发/技改", "研发费用加计、技改、设备补贴、智能制造、数字化转型", "工业/研发类企业适配"],
       ["创新创业载体", "众创空间/孵化器/加速器/产业创新服务平台", "争取载体认定，提升招商可信度"],
       ["三层政策服务", "基础诊断(免费) / 专项申报(收费) / 政府资源对接(深度)", "招商工具 + 盈利工具"]],
      [22, 44, 34])

sheet("09 风险应对", "风险点与应对",
      ["风险", "应对"],
      [["空间合规(用途/消防/承重/用电)", "入驻准入审核；以轻型研发/中试/装配/测试与办公展示为主，拒高污染高噪音高能耗"],
       ["招商去化", "先样板企业 + 首批优惠 + 高校/协会/社群渠道 + 政策抓手；目标务实(年 2000–3000㎡)"],
       ["收益周期", "初期设基础服务费/保底；佣金与去化挂钩；增值服务逐步导入，避免前期重投入"],
       ["政策兑现", "一切以政府最新文件为准；不承诺必得补贴；服务以辅导/申报/对接为主，成功费按实际兑现计"],
       ["月费被复制锁死", "单项目单核算；多项目须重新议价，避免人力成本覆盖不住"]],
      [30, 70])

xlsx_path = os.path.join(HERE, "创智汇6600平合作方案.xlsx")
wb.save(xlsx_path)
print("Excel saved:", xlsx_path, "sheets:", wb.sheetnames)
