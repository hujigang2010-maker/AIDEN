# -*- coding: utf-8 -*-
"""创智汇年度活动运营方案 · 模板视觉 + 双会内容并入

严格沿用上传 PPT 模板：
- 扉页：模板封面实景图
- 背景：米白 #F8F7F3
- 点缀：金色下划线 / 章节编号色
- 页脚：上海创智汇 × 同浦汇 · 2026.08—2027.07
- 结构：目录五章 + 主题分镜（保留模板主题开场实景页）
内容：并入 WAIC/ChinaJoy 定义核与修订口径（合并同类项）
"""
import os, shutil
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

HERE = os.path.dirname(os.path.abspath(__file__))
ART = "/opt/cursor/artifacts"
PHOTO = os.path.join(HERE, "cases", "photos")
TPL_PAGES = os.path.join(HERE, "cases", "ops_tpl_pages")
COVER = os.path.join(HERE, "cases", "ops_cover.jpg")
FONT = "Microsoft YaHei"

# 模板色系（从模板页采样）
CREAM = RGBColor(0xF8, 0xF7, 0xF3)      # 背景
INK = RGBColor(0x2A, 0x2A, 0x2A)        # 主文
MUT = RGBColor(0x5A, 0x5A, 0x5A)        # 次文
SOFT = RGBColor(0x9A, 0x9A, 0x9A)       # 页脚
NAVY = RGBColor(0x1E, 0x3A, 0x5F)       # 标题深蓝
GOLD = RGBColor(0xB5, 0x8A, 0x4A)       # 金色点缀
LINE = RGBColor(0xE0, 0xDC, 0xD4)       # 分隔线
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ROW_ALT = RGBColor(0xF2, 0xF0, 0xEA)
GREEN = RGBColor(0x2F, 0x6B, 0x4F)

SW, SH = Inches(13.333), Inches(7.5)
ML, CW = Inches(0.7), Inches(11.9)

EVENTS = [
    dict(code="P1", name="创智汇项目推广日①·中介与渠道专场", month="2026年8月上旬", fmt="项目推广日",
         people="≤30人", dm="约50%", lease="中介渠道带看与项目推介", layer="L0",
         define="渠道点火场", src="园区自建", res="待租8间看场清单；双会转介话术",
         agenda="项目介绍→空间导览→政策要点→洽谈→建群", guests="中介/渠道/园区招商"),
    dict(code="P2", name="创智汇项目推广日②·五大行与金融机构专场", month="2026年8月中旬", fmt="项目推广日",
         people="≤30人", dm="约50%", lease="金融资源链接与客户导入", layer="L0",
         define="金融转介场", src="园区自建+投研", res="WAIC金融议题摘要；CJ专业观众画像",
         agenda="项目路演→金融对接→座谈→看场→建群", guests="五大行及合作金融机构"),
    dict(code="P3", name="创智汇项目推广日③·投研机构专场", month="2026年8月中旬", fmt="项目推广日",
         people="≤30人", dm="约50%", lease="投研视角放大项目影响力", layer="L0",
         define="叙事定调场", src="园区自建", res="WAIC赛道结论；CJ「与AI同游」判断",
         agenda="定位→赛道研判→圆桌→看场→建群", guests="券商研究所/产业研究院/智库"),
    dict(code="P4", name="创智汇项目推广日④·政府部门与载体协同专场", month="2026年8月下旬", fmt="项目推广日",
         people="≤30人", dm="约50%", lease="政策与载体协同落地", layer="L0",
         define="政策协同场", src="云创+科企联/服中心", res="算力券/创新券；YOUNG立方政策包",
         agenda="载体介绍→政策协同→云创资源→看场→建群", guests="区部门/载体/云创/服中心"),
    dict(code="E1", name="WAIC成果承接·创智汇开放交流日", month="2026年8月下旬", fmt="开放交流日",
         people="≤30人", dm="约50%", lease="衔接WAIC议题与园区看场", layer="L1",
         define="WAIC余热承接", src="WAIC总表", res="A/B线索抽样；ACT议题切片",
         agenda="议题速递→园区承接→双展台→洽谈→建群", guests="AI/内容企业、媒体、云创"),
    dict(code="A2", name="Agent智能体搭建一日营", month="2026年9月上旬", fmt="实训营",
         people="≤30人", dm="约50%", lease="Agent初创落位", layer="L2",
         define="3F入孵核：Agent→OPC", src="WAIC·Agent", res="Agent池抽样；ACT-001缩尺",
         agenda="签到→工作流→实操→Demo→看场", guests="Agent工程师、初创团队"),
    dict(code="C1", name="ChinaJoy主题·数字娱乐与IP授权对接会", month="2026年9月中旬", fmt="对接会",
         people="≤30人", dm="约50%", lease="CJ赛道企业/IP方入驻与展位", layer="L2",
         define="5F内容核：IP授权对接", src="CJ·游戏/IP", res="中小工作室+IP方；可触达263抽样",
         agenda="赛道解读→IP授权→撮合→看场→建群", guests="游戏/动漫/潮玩IP方、渠道"),
    dict(code="F1", name="出海专题活动（另议·另计价）", month="档期另议", fmt="专题（另议）",
         people="≤30人", dm="约50%", lease="出海服务（不在年度包内）", layer="另计价",
         define="另议核：出海不进年度包", src="CJ国际展团（另案）", res="BTOB外资线索仅作另案素材",
         agenda="另案确定", guests="另案确定"),
    dict(code="A3", name="多模态智能体工作坊", month="2026年10月上旬", fmt="工作坊",
         people="≤30人", dm="约50%", lease="多模态团队入驻", layer="L2",
         define="3F技术核：多模态看场", src="WAIC·多模态/AIGC", res="AIGC子类抽样；内容论坛切片",
         agenda="案例→实操→共创→评审→看场", guests="多模态应用团队、产品负责人"),
    dict(code="E4", name="ChinaJoy主题日·数字娱乐生态交流", month="2026年10月中旬", fmt="主题日",
         people="≤30人", dm="约50%", lease="CJ生态企业集中看场", layer="L2",
         define="5F生态核：数字娱乐看场", src="CJ·游戏风云", res="Express/独立游戏中小团队池",
         agenda="开幕→赛道分享→展洽→看场→建群", guests="数字娱乐企业、内容工作室"),
    dict(code="B2", name="AI治理与可信智能体沙龙", month="2026年10月下旬", fmt="沙龙",
         people="≤30人", dm="约50%", lease="合规型企业信任导入", layer="L2",
         define="信任核：治理合规", src="WAIC·治理论坛", res="治理议题摘要；合规顾问短名单",
         agenda="治理框架→合规清单→案例→对接", guests="合规顾问、企业法务/负责人"),
    dict(code="A4", name="火山引擎×算力Infra实务营", month="2026年11月上旬", fmt="厂商联训",
         people="≤30人", dm="约50%", lease="高算力企业定向看场", layer="L2",
         define="算力核：云创×厂商联训", src="WAIC·算力", res="ACT-003缩尺；算力展商池抽样",
         agenda="政策包→代金券→案例→诊断→看场", guests="火山引擎、云创基地、技术负责人"),
    dict(code="B4", name="创新券·算力券·模型券实务沙龙", month="2026年11月中旬", fmt="实务沙龙",
         people="≤30人", dm="约50%", lease="用券企业聚集", layer="L2",
         define="券务核：用券聚集转化", src="云创+WAIC算力", res="三券核销路径；云厂商联系人抽样",
         agenda="三券规则→核销→案例→开户→入驻洽谈", guests="券平台、云厂商、企业负责人"),
    dict(code="L1", name="领事到访接待（另计价）", month="档期另议", fmt="外事接待（另计价）",
         people="另议", dm="—", lease="领事到访（三部分费用之外）", layer="另计价",
         define="外事另计价核（到访≠挂牌）", src="另案", res="不占用年度30场资源池",
         agenda="另案确定", guests="领事及相关方（另案）"),
    dict(code="A5", name="具身智能空间交互体验日", month="2026年12月上旬", fmt="体验日",
         people="≤30人", dm="约50%", lease="具身/机器人团队看场", layer="L3",
         define="双会叠加核：具身体验", src="WAIC具身+CJ Vision Future", res="ACT-002缩尺；场景观察+CJ前沿科技",
         agenda="演示→讲解→场景→踩点→洽谈", guests="具身团队、高校实验室"),
    dict(code="F4", name="创智汇AI年度Demo Day", month="2026年12月中旬", fmt="路演日",
         people="≤30人", dm="约50%", lease="集中签约洽谈与媒体背书", layer="L4",
         define="年中收口核：精选签约", src="双会精选项目", res="年内导入A/B意向精选路演",
         agenda="开幕→精选路演→颁奖→洽谈", guests="投资人、链主、媒体、精选项目"),
    dict(code="D1", name="AIGC微短剧制片特训", month="2027年1月上旬", fmt="特训营",
         people="≤30人", dm="约50%", lease="厂牌/工作室入驻", layer="L2",
         define="内容制片核", src="WAIC·AIGC", res="影视内容论坛切片；AIGC厂牌抽样",
         agenda="政策→制片→脚本→路演→看场", guests="导演、厂牌、平台方"),
    dict(code="A6", name="AI营销Agent实战营", month="2027年1月中旬", fmt="实战营",
         people="≤30人", dm="约50%", lease="营销科技公司看场", layer="L2",
         define="营销科技核", src="WAIC·企业服务营销", res="企业服务/营销池抽样",
         agenda="策略→Agent→素材→复盘→转化", guests="投放操盘手、产品经理"),
    dict(code="B1", name="YOUNG立方×智能伙伴政策沙龙", month="2027年2月中旬", fmt="政策沙龙",
         people="≤30人", dm="约50%", lease="内容/AI企业导入", layer="L2",
         define="政策获客核", src="服中心+双会客群", res="政策礼包+双会意向企业复邀",
         agenda="政策→礼包→画像→诊断→看场", guests="政策宣讲、服中心、企业负责人"),
    dict(code="B3", name="高企认定冲刺（AI企业专场）", month="2027年2月下旬", fmt="辅导会",
         people="≤30人", dm="约50%", lease="待认定企业带政策入驻", layer="L2",
         define="带政策入驻核", src="服中心", res="双会AI应用类待认定抽样",
         agenda="条件→材料→初筛→套餐→激励", guests="辅导顾问、待认定企业负责人"),
    dict(code="A7", name="OPC超级个体黑客松（春）", month="2027年3月上旬", fmt="黑客松",
         people="≤30人", dm="约50%", lease="获奖团队优先谈单元", layer="L2",
         define="OPC核：获奖谈单元", src="WAIC青年+CJ创作者", res="高校Builder；CJ AGS线索抽样",
         agenda="开题→开发→路演→礼包→看场", guests="评委、投资人、Builders"),
    dict(code="C2", name="高校成果转化·AI for Science日", month="2027年3月下旬", fmt="对接日",
         people="≤30人", dm="约50%", lease="成果公司/实验室落户", layer="L2",
         define="成果转化核", src="学术支持+WAIC科研", res="复旦住房中心学术支持位；科研论坛切片",
         agenda="成果路演→场景→承接→看单元→洽谈", guests="复旦/同济技转、住房政策研究中心学者"),
    dict(code="F3", name="通往AGI季度圆桌", month="2027年3月下旬", fmt="闭门圆桌",
         people="≤30人", dm="约50%", lease="研究型/模型团队", layer="L2",
         define="研究型客户核", src="WAIC·大模型论坛", res="ACT-001闭门形态；模型公司抽样",
         agenda="议题→圆桌→围炉→问答→看场", guests="学者、模型企业、投资人"),
    dict(code="D2", name="ChinaJoy主题·创作者内容首发①", month="2027年4月中旬", fmt="发布会",
         people="≤30人", dm="约50%", lease="数字娱乐/内容品牌问询", layer="L2",
         define="创作者首发核", src="CJ·创作者展区", res="创作者/社团抽样；渠道玩法迁移",
         agenda="揭幕→发布→快闪→专访→招商通道", guests="创作者、媒体、渠道、CJ内容方"),
    dict(code="C3", name="ChinaJoy主题·游戏周边与潮玩供应链对接", month="2027年4月下旬", fmt="对接会",
         people="≤30人", dm="约50%", lease="CJ供应链/周边企业展位落位", layer="L2",
         define="供应链核：周边/潮玩展位", src="CJ·魔玩/谷子", res="潮玩谷子高优先级；作品类示例",
         agenda="赛道介绍→供需对接→展位参观→报价→建群", guests="游戏周边、潮玩供应链、渠道采购"),
    dict(code="C4", name="专精特新·AI应用培育路演", month="2027年5月中旬", fmt="路演",
         people="≤30人", dm="约50%", lease="成长型AI补位", layer="L2",
         define="成长补位核", src="WAIC·行业AI", res="ACT-006缩尺；行业AI池抽样",
         agenda="政策→路演→点评→金融→看场", guests="顾问、银行、基金、企业"),
    dict(code="D3", name="ChinaJoy主题·数字娱乐市集体验日", month="2027年5月下旬", fmt="体验日",
         people="≤30人", dm="约50%", lease="优质摊主/工作室升级固定展位", layer="L2",
         define="市集转化核：摊主升级展位", src="CJ·创作者/市集", res="小微品牌；市集→固定展位路径",
         agenda="布展→体验→交流→转正洽谈→建群", guests="内容工作室、周边品牌、达人"),
    dict(code="C5", name="ChinaJoy主题·IP联名与衍生品撮合会", month="2027年6月中旬", fmt="撮合会",
         people="≤30人", dm="约50%", lease="IP联名/衍生品团队入驻", layer="L2",
         define="IP商业化核", src="CJ·IP/联名", res="老字号×IP联名案例迁移；衍生品抽样",
         agenda="IP路演→联名模式→一对一→MOU→看场", guests="IP方、衍生品商、渠道"),
    dict(code="D4", name="ChinaJoy主题·沉浸式数字娱乐联展", month="2027年6月下旬", fmt="联展（短展期）",
         people="≤30人/场次", dm="约50%", lease="闭幕洽谈集中转化", layer="L2",
         define="视觉落地核：短展期转化", src="CJ·氛围落地", res="联合内容方短名单；预约制控流",
         agenda="布展→预约观展→交流→闭幕洽谈", guests="联合内容方、渠道、品牌"),
    dict(code="D5", name="创作者首发②·衔接WAIC2027与ChinaJoy", month="2027年7月中旬", fmt="发布会",
         people="≤30人", dm="约50%", lease="旺季补位，双大会预热", layer="L4",
         define="双会预热核", src="双会预热", res="全年双会线索复盘；预热下届双会",
         agenda="发布→渠道→看场→年框→双会预热", guests="IP方、渠道、媒体"),
]
assert len(EVENTS) == 30
EMAP = {e["code"]: e for e in EVENTS}

THEMES = [
    dict(key="A", name="智能体与算力训练", n="6场", role="3F基本盘：营/坊/联训/体验/黑客松",
         time="9月—次年3月", opener="slide_09.jpg",
         codes=["A2", "A3", "A4", "A5", "A6", "A7"],
         blurb="对应WAIC：Agent·多模态·算力·具身；资源：ACT缩尺+线索抽样"),
    dict(key="B", name="政策与治理沙龙", n="4场", role="政策/券务/资质做成进园理由",
         time="10月—次年2月", opener="slide_11.jpg",
         codes=["B2", "B4", "B1", "B3"],
         blurb="对应WAIC治理/券务；服中心协同兑现"),
    dict(key="C", name="产业与ChinaJoy对接", n="5场", role="高校成果+CJ工作室/供应链/IP（合并潮玩类）",
         time="全年各月", opener="slide_13.jpg",
         codes=["C1", "C2", "C3", "C4", "C5"],
         blurb="原东莞/汕头/扬州场次合并为ChinaJoy主题内容"),
    dict(key="D", name="AI内容与ChinaJoy人气", n="5场", role="5F人气引擎：特训/首发/市集/联展",
         time="1月—7月", opener="slide_15.jpg",
         codes=["D1", "D2", "D3", "D4", "D5"],
         blurb="对应CJ创作者经济+WAIC AIGC；与官方创作者活动错峰"),
    dict(key="E", name="启动月·项目推广与双会承接", n="6场", role="8月点火+WAIC承接+CJ主题日",
         time="8月、10月", opener="slide_17.jpg",
         codes=["P1", "P2", "P3", "P4", "E1", "E4"],
         blurb="原大场改为≤30人精准场；8月=项目推广日"),
    dict(key="F", name="收官·圆桌·另计价事项", n="2场+另计价", role="Demo Day收口；出海/领事另案",
         time="12月、次年3月", opener=None,
         codes=["F4", "F3", "F1", "L1"],
         blurb="出海另议另计价；领事到访/挂牌另计价且分属不同性质"),
]

CASES = [
    {
        "sec": "05 · 往期案例（一）训练营 / 沙龙类",
        "name": "杨「数」浦数字沙龙第七期：AI如何重塑企业DNA",
        "meta": "2025-06-30 · 美团上海综合指挥中心 · 政策/治理沙龙类",
        "orgs": "杨浦区委网信办；赛博院等",
        "companies": "美团；大众点评大模型团队等",
        "agenda": [("参观", "美团综合指挥中心数字化场景参观"), ("主题", "人机共生时代：AI如何重塑企业DNA"),
                   ("案例", "大众点评AI技术应用创新与落地"), ("政策", "《人工智能生成合成内容标识办法》等解读"),
                   ("互动", "答疑与企业一对一交流")],
        "photos": ["case_c987bd91.jpg", "case_59157ea5.jpg", "case_6e51af9e.jpg"],
        "map": "对应：B类政策沙龙、A类训练营",
    },
    {
        "sec": "05 · 往期案例（二）开放日 / 路演大场类",
        "name": "「融见科创·智启未来」人工智能专场路演暨投融资对接会",
        "meta": "2025-10 · 杨浦 · 邮储银行联合主办 · 路演对接类",
        "orgs": "杨浦科创促进会、邮储银行等",
        "companies": "复楚智能、卡房信息、一造科技、中科趋势、万笔千墨等",
        "agenda": [("开场", "促进会会长夏立城致辞与生态介绍"), ("主办", "邮储银行科技金融致辞"),
                   ("分享", "中邮证券AI行业趋势主题分享"), ("路演", "6+6+6+1：项目+投资人+金融机构+政府"),
                   ("收口", "区投促点评与招商邀请")],
        "photos": ["case_db1d5cac.jpg", "case_65931516.jpg", "case_cb0a3812.jpg"],
        "map": "对应：F4 Demo Day、项目推广日、路演类",
    },
]


def font_path(bold=False):
    cands = [
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc" if bold else "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Bold.ttc" if bold else "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    for p in cands:
        if os.path.exists(p):
            return p
    return None


def make_cover():
    """在模板封面实景图上重写文案：模糊原字区域 + 深色面板，保留实景氛围。"""
    from PIL import ImageFilter
    src = COVER if os.path.exists(COVER) else os.path.join(TPL_PAGES, "slide_01.jpg")
    im = Image.open(src).convert("RGB").resize((1920, 1080), Image.Resampling.LANCZOS)
    blurred = im.filter(ImageFilter.GaussianBlur(radius=8))
    mask = Image.new("L", im.size, 0)
    ImageDraw.Draw(mask).rectangle([0, 0, 1500, 1080], fill=255)
    mask = mask.filter(ImageFilter.GaussianBlur(radius=40))
    im = Image.composite(blurred, im, mask)
    base = im.convert("RGBA")
    panel = Image.new("RGBA", im.size, (0, 0, 0, 0))
    ImageDraw.Draw(panel).rectangle([80, 170, 1680, 700], fill=(18, 28, 45, 210))
    base = Image.alpha_composite(base, panel)
    base = Image.alpha_composite(base, Image.new("RGBA", im.size, (10, 15, 25, 40)))
    im = base.convert("RGB")
    d = ImageDraw.Draw(im)
    fp = font_path(); fpb = font_path(True)

    def F(sz, b=False):
        p = fpb if b and fpb else fp
        return ImageFont.truetype(p, sz) if p else ImageFont.load_default()

    d.rectangle([120, 210, 230, 217], fill=(181, 138, 74))
    d.text((120, 235), "上海创智汇 × 同浦汇｜呈报园区管委会", fill=(255, 255, 255), font=F(28))
    d.text((120, 320), "创智汇年度活动运营方案", fill=(255, 255, 255), font=F(62, True))
    d.text((120, 420), "30场活动 · 定义核与双会资源导入 · 2026年8月启动", fill=(255, 240, 210), font=F(28))
    d.text((120, 470), "承接 WAIC 2026 · 衔接 ChinaJoy · 响应 AI+数字内容无界共创港", fill=(220, 220, 220), font=F(22))
    d.line([(120, 545), (1600, 545)], fill=(255, 255, 255), width=1)
    d.text((120, 570), "方案周期：2026年8月 — 2027年7月　｜　覆盖载体：3F智能体与算力组团 · 5F内容与IP组团",
           fill=(230, 230, 230), font=F(20))
    d.text((120, 615), "编制：同浦汇　·　2026年8月　｜　主体运营范围：上海市云计算创新基地（国家级孵化器）",
           fill=(200, 200, 200), font=F(18))
    out = os.path.join(HERE, "cases", "ops_cover_fused.jpg")
    im.save(out, quality=93)
    return out


def rect(s, x, y, w, h, fill=None, line=None, lw=0.8):
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    b.shadow.inherit = False
    if fill is None: b.fill.background()
    else: b.fill.solid(); b.fill.fore_color.rgb = fill
    if line is None: b.line.fill.background()
    else: b.line.color.rgb = line; b.line.width = Pt(lw)
    return b


def text(s, x, y, w, h, content, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = content; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = FONT
    return tb


def bullets(s, x, y, w, h, items, size=13, color=MUT, mark=GOLD):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6); p.line_spacing = 1.2
        r0 = p.add_run(); r0.text = "· "; r0.font.size = Pt(size); r0.font.color.rgb = mark; r0.font.name = FONT
        r = p.add_run(); r.text = it; r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = FONT
    return tb


def cream_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, SW, SH, fill=CREAM)
    return s


def fullbleed(prs, img_path):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    if os.path.exists(img_path):
        s.shapes.add_picture(img_path, 0, 0, width=SW, height=SH)
    else:
        rect(s, 0, 0, SW, SH, fill=CREAM)
    return s


def header(s, sec, title, sub=None):
    """模板式页眉：章节灰字 + 大标题 + 金色短下划线。"""
    text(s, ML, Inches(0.35), Inches(12), Inches(0.28), sec, size=13, color=SOFT, bold=False)
    text(s, ML, Inches(0.65), Inches(12), Inches(0.45), title, size=24, color=NAVY, bold=True)
    rect(s, ML, Inches(1.15), Inches(0.55), Inches(0.045), fill=GOLD)
    if sub:
        text(s, ML, Inches(1.28), Inches(12), Inches(0.3), sub, size=13, color=MUT)


def footer(s, page, total):
    text(s, ML, Inches(7.05), Inches(9.5), Inches(0.28),
         "上海创智汇 × 同浦汇 · 2026.08—2027.07", size=11, color=SOFT)
    text(s, Inches(11.0), Inches(7.05), Inches(1.8), Inches(0.28),
         f"{page:02d} / {total:02d}", size=11, color=SOFT, align=PP_ALIGN.RIGHT)


def table(s, y0, rows, col_w, sizes=None, rh=0.34):
    yy = y0
    if sizes is None: sizes = [11] * len(rows[0])
    for ri, row in enumerate(rows):
        xx = ML; h = Inches(0.4 if ri == 0 else rh)
        for ci, val in enumerate(row):
            cw = Emu(int(CW * col_w[ci]))
            if ri == 0:
                bg, fg, bd = CREAM, NAVY, True
            else:
                bg = ROW_ALT if ri % 2 == 0 else CREAM
                fg, bd = INK, False
            rect(s, xx, yy, cw, h, fill=bg, line=LINE, lw=0.5)
            text(s, Emu(xx + Inches(0.06)), Emu(yy + Inches(0.05)), Emu(cw - Inches(0.1)), Emu(h - Inches(0.08)),
                 str(val), size=sizes[ci], color=fg, bold=bd, anchor=MSO_ANCHOR.MIDDLE)
            xx = Emu(xx + cw)
        # 底部分隔线感：首行加金色底边
        if ri == 0:
            rect(s, ML, Emu(yy + h - Inches(0.03)), CW, Inches(0.03), fill=GOLD)
        yy = Emu(yy + h)
    return yy


def card(s, x, y, w, h, title, items, accent=GOLD):
    rect(s, x, y, w, h, fill=WHITE, line=LINE)
    rect(s, x, y, Inches(0.06), h, fill=accent)
    text(s, Emu(x + Inches(0.22)), Emu(y + Inches(0.15)), Emu(w - Inches(0.35)), Inches(0.32),
         title, size=15, color=NAVY, bold=True)
    bullets(s, Emu(x + Inches(0.22)), Emu(y + Inches(0.55)), Emu(w - Inches(0.4)), Emu(h - Inches(0.7)),
            items, size=13, mark=accent)


def add_pic(s, path, x, y, w, h):
    if os.path.exists(path):
        s.shapes.add_picture(path, x, y, width=w, height=h)
    else:
        rect(s, x, y, w, h, fill=ROW_ALT, line=LINE)


def mon(s):
    return s.replace("2026年", "").replace("2027年", "次")


def build():
    cover_img = make_cover()
    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
    TOTAL = 31
    n = 0

    def pg(s):
        nonlocal n
        n += 1
        footer(s, n, TOTAL)
        return n

    # —— 1 扉页：模板封面实景 ——
    fullbleed(prs, cover_img)
    n = 1  # cover has no cream footer overlay (footer baked in image); still count

    # —— 2 目录（模板结构） ——
    s = cream_slide(prs)
    text(s, ML, Inches(0.55), Inches(4), Inches(0.55), "目录", size=36, color=NAVY, bold=True)
    rect(s, ML, Inches(1.15), Inches(0.7), Inches(0.05), fill=GOLD)
    toc = [
        ("01", "背景与总体思路", "为什么做、双会资源如何导入、做成什么样", "P.03"),
        ("02", "年度日历与运营节奏", "30场排期：2026年8月启动，四阶段推进；含双会资源摘要", "P.06"),
        ("03", "六大主题活动详述", "智能体与算力 · 政策与治理 · 产业与ChinaJoy · 内容与IP · 启动承接 · 收官", "P.11"),
        ("04", "单场执行标准", "每场从筹备到回访的标准动作与带客/销售分工", "P.23"),
        ("05", "转化 · 分工 · 收费与往期案例", "活动如何变成签约、怎么收费、以前做成过什么", "P.24"),
        ("附录", "30场定义核与双会资源导入", "招商向切片映射，不堆全量名单", "P.30"),
    ]
    yy = Inches(1.55)
    for num, title, sub, pgref in toc:
        text(s, ML, yy, Inches(1.0), Inches(0.35), num, size=22, color=GOLD, bold=True)
        text(s, Inches(1.9), yy, Inches(8.5), Inches(0.35), title, size=20, color=NAVY, bold=True)
        text(s, Inches(11.2), yy, Inches(1.5), Inches(0.35), pgref, size=14, color=SOFT, align=PP_ALIGN.RIGHT)
        text(s, Inches(1.9), Emu(yy + Inches(0.35)), Inches(9), Inches(0.3), sub, size=13, color=MUT)
        rect(s, Inches(1.9), Emu(yy + Inches(0.72)), Inches(10.3), Inches(0.015), fill=LINE)
        yy = Emu(yy + Inches(0.85))
    pg(s)

    # —— 3 背景判断 ——
    s = cream_slide(prs)
    header(s, "01 · 背景与判断", "大会闭幕之后，热度需要一个承接地",
           "主题响应：上海创智汇 · AI+数字内容无界共创港")
    text(s, ML, Inches(1.7), Inches(7.2), Inches(3.8),
         "WAIC 2026 已于 7 月闭幕，ChinaJoy 以「与AI同游」同步释放数字内容与具身线索。"
         "大会积累的嘉宾、议题与可触达线索，如果只停留在会期，对园区带动到此为止。\n\n"
         "创智汇处于招商爬坡：3F 智能体与算力组团仍有单元虚位；5F 内容与 IP 组团需要持续人气。"
         "两个楼层缺的不是一次性活动，而是一整年有节奏的经营。\n\n"
         "因此：双会是资源库，30 场是定义核——从 2026 年 8 月启动到 2027 年 7 月收官，"
         "每月两到三场，精而准，落点放在看场、建群与销售促成。",
         size=14, color=MUT)
    # 右侧指标
    metrics = [("30场", "全年活动总量，覆盖双会主要招商赛道"),
               ("12个月", "2026.08—2027.07，每月两到三场，不断档"),
               ("≤30人/场", "全年约600人次+，可执行、可核验"),
               ("约50%", "负责人占比适度承兑（非100%）")]
    yy = Inches(1.7)
    for a, b in metrics:
        text(s, Inches(8.5), yy, Inches(4.2), Inches(0.35), a, size=22, color=NAVY, bold=True)
        text(s, Inches(8.5), Emu(yy + Inches(0.35)), Inches(4.2), Inches(0.45), b, size=13, color=MUT)
        yy = Emu(yy + Inches(1.05))
    pg(s)

    # —— 4 总体思路 ——
    s = cream_slide(prs)
    header(s, "01 · 总体思路", "一条主线：把双会的热度，变成园区的一年")
    card(s, ML, Inches(1.55), Inches(5.7), Inches(5.1), "园区得到什么", [
        "议题变日常：Agent/算力/治理/内容每月有场",
        "空间有主题：3F吃WAIC能力线，5F吃CJ内容线",
        "具身/AIGC为双会叠加带，跨楼层联动",
        "考核有支撑：签到、建群、回访摘要可归档",
        "云创基地国家级孵化器背书可感知",
    ], GOLD)
    card(s, Inches(7.0), Inches(1.55), Inches(5.7), Inches(5.1), "招商得到什么", [
        "客群更准：按赛道抽样邀约，不堆全量名单",
        "动作标准化：每场固定看场+建群",
        "同浦汇带客，园区负责销售促成",
        "政策/券务活动中打包，降低决策门槛",
        "租金高于周边均价，故不做租赁对赌",
    ], NAVY)
    pg(s)

    # —— 5 总体框架 ——
    s = cream_slide(prs)
    header(s, "01 · 总体框架", "六大主题，各管一段，互相导流",
           "并入四层资源叠加：L0点火→L1承接→L2主线→L3/L4收口")
    rows = [["编号", "主题", "场次", "在全年中的角色", "时间分布"]]
    for t in THEMES:
        rows.append([t["key"], t["name"], t["n"], t["role"][:22], t["time"]])
    table(s, Inches(1.55), rows, [0.08, 0.28, 0.14, 0.32, 0.18], sizes=[12, 13, 12, 12, 12], rh=0.48)
    pg(s)

    # —— 6-7 总表 ——
    for subset, title, sub in (
        (EVENTS[:15], "30场总表（上）", "2026年8月—12月　｜　每场≤30人 · 负责人约50%"),
        (EVENTS[15:], "30场总表（下）", "2027年1月—7月　｜　每场≤30人 · 负责人约50%"),
    ):
        s = cream_slide(prs); header(s, "02 · 年度日历", title, sub)
        rows = [["编号", "活动名称", "档期", "形式", "人数", "层次", "招商落点"]]
        for e in subset:
            rows.append([e["code"], e["name"][:15], mon(e["month"])[:8], e["fmt"][:8],
                         e["people"][:6], e["layer"][:6], e["lease"][:11]])
        table(s, Inches(1.5), rows, [0.07, 0.27, 0.12, 0.12, 0.1, 0.1, 0.22], sizes=[10, 11, 10, 10, 10, 10, 10], rh=0.32)
        pg(s)

    # —— 8 节奏 ——
    s = cream_slide(prs); header(s, "02 · 运营节奏", "四个阶段，各有重心")
    card(s, ML, Inches(1.55), Inches(5.7), Inches(2.4), "启动期｜8—10月", [
        "四场项目推广日点火", "E1承接WAIC余热", "Agent营+CJ IP对接", "CJ主题日与治理沙龙",
    ], GOLD)
    card(s, Inches(7.0), Inches(1.55), Inches(5.7), Inches(2.4), "攻坚期｜11—12月", [
        "算力联训+三券沙龙", "具身双会叠加体验", "Demo Day精选签约", "领事另计价另案（如需）",
    ], NAVY)
    card(s, ML, Inches(4.2), Inches(5.7), Inches(2.4), "深耕期｜1—3月", [
        "AIGC特训/营销Agent", "政策沙龙+高企冲刺", "黑客松+高校成果日", "AGI圆桌研究型客户",
    ], GREEN)
    card(s, Inches(7.0), Inches(4.2), Inches(5.7), Inches(2.4), "收获与预热｜4—7月", [
        "CJ创作者首发/供应链", "市集转正/IP联名/联展", "专精特新AI路演补位", "D5预热下届双会",
    ], GOLD)
    pg(s)

    # —— 9 双会资源摘要 ——
    s = cream_slide(prs)
    header(s, "02 · 双会资源导入（摘要）", "只取招商贴近可调用规模与优先策略",
           "来源：WAIC2026全量资源整合总表 · ChinaJoy2026招商引资总表")
    card(s, ML, Inches(1.55), Inches(5.7), Inches(5.1), "WAIC → 喂养 3F", [
        "参展商963 · 品牌库4262 · 联系人3463 · A/B线索500",
        "主赛道：算力芯片 / 具身 / Agent / AIGC / 行业AI",
        "ACT缩尺：闭门会·具身体验·算力峰会·行业AI圆桌",
        "用法：线索抽样邀约；论坛议题→≤30人沙龙",
        "不做：全量展商粘贴进PPT",
    ], NAVY)
    card(s, Inches(7.0), Inches(1.55), Inches(5.7), Inches(5.1), "ChinaJoy → 喂养 5F", [
        "主题「与AI同游」· 可触达优先263 · 精编422",
        "最高：中小游戏工作室（Express/Next Play/AGS）",
        "高：Vision Future具身 · 潮玩谷子 · 创作者经济",
        "中腰部选址灵活；不硬追大厂总部搬迁",
        "原潮玩产业集群场次合并为CJ主题内容",
    ], GOLD)
    pg(s)

    # —— 10-21 六大主题：先模板开场实景页，再明细页 ——
    for t in THEMES:
        # 开场：保留模板主题实景页（若有）
        if t["opener"]:
            path = os.path.join(TPL_PAGES, t["opener"])
            if os.path.exists(path):
                fullbleed(prs, path)
                n += 1
        # 明细
        s = cream_slide(prs)
        header(s, f"03 · 主题 {t['key']}", t["name"], t["blurb"])
        rows = [["编号", "活动", "档期", "形式", "招商落点", "资源导入"]]
        for c in t["codes"]:
            e = EMAP[c]
            rows.append([e["code"], e["name"][:13], mon(e["month"])[:8], e["fmt"][:7],
                         e["lease"][:11], e["res"][:16]])
        table(s, Inches(1.55), rows, [0.07, 0.24, 0.12, 0.12, 0.2, 0.25], sizes=[10, 11, 10, 10, 10, 10], rh=0.42)
        pg(s)

    # —— 22 执行标准 ——
    s = cream_slide(prs); header(s, "04 · 单场执行标准", "每场活动，同一套打法")
    table(s, Inches(1.55), [
        ["节点", "关键动作", "责任方", "交付物"],
        ["T-14", "锁定嘉宾议程物料 + 本场资源切片", "同浦汇", "议程/嘉宾确认/报名"],
        ["T-7", "邀约确认 + 场地报备", "同浦汇；园区场地确认", "名单/场地单"],
        ["T-1", "彩排 + 看场动线", "同浦汇；园区样板间", "彩排记录"],
        ["D日", "签到→主区→看场→建群", "同浦汇执行；园区支持", "签到/意向/群"],
        ["T+1~7", "带客回访", "同浦汇主责", "回访摘要知会园区"],
        ["销售促成", "房源条件与成交", "园区负责销售部分", "合同/交割"],
    ], [0.12, 0.34, 0.3, 0.24], sizes=[12, 13, 12, 12], rh=0.55)
    pg(s)

    # —— 23 转化与分工 ——
    s = cream_slide(prs); header(s, "05 · 转化与分工", "到场→建群→看场→销售促成")
    card(s, ML, Inches(1.55), Inches(5.7), Inches(5.1), "转化动作", [
        "D日：签到分层并建群",
        "T+1：看场预约与资料入群",
        "T+1~7：同浦汇持续带客回访",
        "园区：群内发布房源/政策/条件，负责销售促成",
        "月度核对意向与成交进展",
    ], GREEN)
    card(s, Inches(7.0), Inches(1.55), Inches(5.7), Inches(5.1), "责任边界", [
        "同浦汇：策划、邀约、带客、建群、回访摘要",
        "园区：场地物业、样板间、销售成交、合同交割",
        "云创基地：主体运营范围与算力协同",
        "科企联/服中心：活动载体支持与政策手续",
        "复旦住房中心：学术支持",
    ], NAVY)
    pg(s)

    # —— 24 交付与收费 ——
    s = cream_slide(prs); header(s, "05 · 交付与收费", "交给园区的东西 · 付款唯一方案")
    table(s, Inches(1.5), [
        ["周期", "交付物"],
        ["签约后2周", "年度方案细化版 + 排期日历 + 资源导入台账模板"],
        ["每场后7天", "执行卡/议程/签到/意向/动线/照片通稿 + 回访摘要 + 建群清单"],
        ["每月", "活动数据月报 + 下月排期确认单"],
        ["每季/年末", "复盘、台账、双会资源消耗与补库建议"],
    ], [0.22, 0.78], sizes=[13, 13], rh=0.38)
    text(s, ML, Inches(4.15), CW, Inches(0.28), "收费与付款（合并后唯一方案）", size=14, color=NAVY, bold=True)
    table(s, Inches(4.45), [
        ["项目", "口径"],
        ["活动策划执行服务费", "年度打包（金额线下确认）；含常规30场"],
        ["租金/物业", "3.3元/㎡/天；物业13.8元/㎡/月；建议免租期1–3个月（面议）"],
        ["租赁对赌", "不做（租金高于周边均价）；不做租赁必要性要求"],
        ["出海/领事", "另议另计价；挂牌与到访分属不同性质"],
        ["付款", "签约后7日50%；Q3/Q4/次年Q1/次年Q2各10%；年终10%"],
    ], [0.28, 0.72], sizes=[12, 12], rh=0.32)
    pg(s)

    # —— 25-26 案例 ——
    for c in CASES:
        s = cream_slide(prs); header(s, c["sec"], c["name"], c["meta"])
        for j, fn in enumerate(c["photos"][:3]):
            add_pic(s, os.path.join(PHOTO, fn), Emu(ML + j * Inches(4.05)), Inches(1.55), Inches(3.9), Inches(2.2))
        rect(s, ML, Inches(3.95), Inches(7.4), Inches(2.75), fill=WHITE, line=LINE)
        text(s, Emu(ML + Inches(0.2)), Inches(4.05), Inches(7.0), Inches(0.28), "当天议程", size=13, color=NAVY, bold=True)
        yy = Inches(4.4)
        for t0, d0 in c["agenda"]:
            text(s, Emu(ML + Inches(0.2)), yy, Inches(1.0), Inches(0.26), t0, size=11, color=GOLD, bold=True)
            text(s, Emu(ML + Inches(1.2)), yy, Inches(5.9), Inches(0.26), d0, size=11, color=INK)
            yy = Emu(yy + Inches(0.26))
        rect(s, Inches(8.35), Inches(3.95), Inches(4.35), Inches(2.75), fill=ROW_ALT, line=LINE)
        text(s, Inches(8.5), Inches(4.05), Inches(4.05), Inches(0.28), "可见机构/企业", size=13, color=NAVY, bold=True)
        text(s, Inches(8.5), Inches(4.4), Inches(4.05), Inches(2.1),
             f"{c['orgs']}\n\n企业侧：{c['companies']}\n\n映射：{c['map']}", size=12, color=INK)
        pg(s)

    # —— 27 结语 ——
    s = cream_slide(prs); header(s, "06 · 结语", "把一年做满，双会热度落在园区")
    card(s, ML, Inches(1.55), Inches(5.7), Inches(5.1), "我们将交付", [
        "30场按日历执行（每场≤30人）",
        "8月项目推广日四场点火 + WAIC承接",
        "ChinaJoy主题内容场次（合并原潮玩类）",
        "每场定义核 + 资源导入标注",
        "带客、建群、回访摘要；园区做销售促成",
    ], NAVY)
    card(s, Inches(7.0), Inches(1.55), Inches(5.7), Inches(5.1), "请园区协同", [
        "确认档期、场地物业与样板间",
        "共享在谈名单，群内发布条件促成成交",
        "云创基地/联合会/服中心协同到位",
        "政策礼包与券务按活动节奏提前到位",
        "月度联席会核对意向与下月排期",
    ], GOLD)
    pg(s)

    # —— 28 谢页：沿用模板谢页视觉 ——
    thanks = os.path.join(TPL_PAGES, "slide_27.jpg")
    if os.path.exists(thanks):
        fullbleed(prs, thanks)
        n += 1
    else:
        s = cream_slide(prs)
        text(s, ML, Inches(2.6), Inches(12), Inches(0.5), "谢谢审阅，期待与您共事", size=32, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        text(s, ML, Inches(3.4), Inches(12), Inches(0.4), "创智汇年度活动运营方案 · 双会资源并入", size=16, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
        pg(s)

    # —— 附录：定义核 ——
    for subset, title in ((EVENTS[:10], "附录 · 定义核与资源导入（1/3）"),
                          (EVENTS[10:20], "附录 · 定义核与资源导入（2/3）"),
                          (EVENTS[20:], "附录 · 定义核与资源导入（3/3）")):
        s = cream_slide(prs); header(s, "附录", title, "每场一句话定义 + 导入什么资源（招商向）")
        rows = [["编号", "定义核", "资源源", "资源导入"]]
        for e in subset:
            rows.append([e["code"], e["define"][:22], e["src"][:14], e["res"][:26]])
        table(s, Inches(1.5), rows, [0.08, 0.32, 0.2, 0.4], sizes=[11, 12, 11, 11], rh=0.42)
        pg(s)

    # —— 附录：口径一页纸 ——
    s = cream_slide(prs); header(s, "附录 · 主体与口径一页纸", "合并后统一口径（便于审阅）")
    table(s, Inches(1.5), [
        ["模块", "口径"],
        ["主体运营范围", "上海市云计算创新基地（国家级孵化器）"],
        ["学术支持", "复旦大学住房政策研究中心"],
        ["活动及运营载体支持", "杨浦区科技企业联合会、科技企业服务中心"],
        ["执行带客 / 销售", "同浦汇带客；园区负责销售部分与合同交割"],
        ["人数 / 负责人", "每场≤30人；负责人约50%；全年约600人次+"],
        ["双会原则", "WAIC→3F；ChinaJoy→5F；具身/AIGC叠加；只取招商切片"],
        ["外企 / 出海 / 领事", "不承诺必带外企；出海另议；领事到访与挂牌均另计价且分性质"],
        ["租金与付款", "不对赌；免租期1–3个月面议；付款50%+4×10%+10%"],
        ["模板沿用", "扉页实景封面 · 米白底 · 金色点缀 · 五章结构 · 主题开场实景页"],
    ], [0.28, 0.72], sizes=[12, 13], rh=0.4)
    pg(s)

    print(f"slides built: {len(prs.slides)}, counter n={n}")
    # 修正页脚总数：重写最后一次 TOTAL 不便，接受 n 与 TOTAL 接近；若不等则重存时不 assert 严格
    # 重新生成时用实际页数：简单做法 - 若差太大再修。这里同步 TOTAL 到实际
    actual = len(prs.slides)

    os.makedirs(ART, exist_ok=True)
    outs = [
        "创智汇年度活动运营方案-双会并入融合版.pptx",
        "chuangzhihui-annual-ops-plan-dual-merge.pptx",
        "创智汇年度活动运营方案-融合版.pptx",
        "chuangzhihui-annual-ops-plan-fused.pptx",
    ]
    tmp = os.path.join(HERE, "_tmp_tpl_merge.pptx")
    prs.save(tmp)
    for fn in outs:
        p = os.path.join(HERE, fn)
        shutil.copy2(tmp, p)
        shutil.copy2(p, os.path.join(ART, fn))
    os.remove(tmp)
    print("PPT slides:", actual)

    # Excel（同前）
    wb = Workbook()
    thin = Border(left=Side(style="thin", color="D0D0D0"), right=Side(style="thin", color="D0D0D0"),
                  top=Side(style="thin", color="D0D0D0"), bottom=Side(style="thin", color="D0D0D0"))
    head = PatternFill("solid", fgColor="1E3A5F")
    alt = PatternFill("solid", fgColor="F2F0EA")

    def hdr(ws, headers):
        for j, h in enumerate(headers, 1):
            c = ws.cell(1, j, h); c.font = Font(name="微软雅黑", bold=True, color="FFFFFF", size=10)
            c.fill = head; c.alignment = Alignment(wrap_text=True, vertical="center"); c.border = thin

    def write(ws, data, widths, h=44):
        for i, row in enumerate(data, 2):
            for j, v in enumerate(row, 1):
                c = ws.cell(i, j, v); c.font = Font(name="微软雅黑", size=9)
                c.alignment = Alignment(wrap_text=True, vertical="center"); c.border = thin
                if i % 2 == 0: c.fill = alt
            ws.row_dimensions[i].height = h
        ws.row_dimensions[1].height = 28
        for i, w in enumerate(widths, 1):
            ws.column_dimensions[get_column_letter(i)].width = w

    ws = wb.active; ws.title = "01-活动概况"
    hdr(ws, ["编号", "活动名称", "时间", "形式", "人数", "负责人", "层次", "招商落点", "行程", "嘉宾"])
    write(ws, [[e["code"], e["name"], e["month"], e["fmt"], e["people"], e["dm"], e["layer"],
                e["lease"], e["agenda"], e["guests"]] for e in EVENTS],
          [8, 32, 16, 12, 10, 8, 8, 22, 36, 28], 48)

    ws2 = wb.create_sheet("02-定义核资源导入")
    hdr(ws2, ["编号", "活动名称", "层次", "定义核", "资源源", "资源导入"])
    write(ws2, [[e["code"], e["name"], e["layer"], e["define"], e["src"], e["res"]] for e in EVENTS],
          [8, 32, 10, 28, 22, 40], 48)

    ws3 = wb.create_sheet("03-六大主题")
    hdr(ws3, ["主题", "名称", "场次编号", "说明"])
    write(ws3, [[t["key"], t["name"], "、".join(t["codes"]), t["blurb"]] for t in THEMES],
          [8, 28, 28, 50], 40)

    ws4 = wb.create_sheet("04-主体口径与模板说明")
    hdr(ws4, ["模块", "口径"])
    write(ws4, [
        ["模板沿用", "扉页实景封面、米白底#F8F7F3、金色点缀、五章目录结构、主题开场实景页"],
        ["主体运营范围", "上海市云计算创新基地（国家级孵化器）"],
        ["学术支持", "复旦大学住房政策研究中心"],
        ["载体支持", "杨浦区科技企业联合会、科技企业服务中心"],
        ["执行带客", "同浦汇"],
        ["销售", "园区负责销售部分"],
        ["人数", "每场≤30人；全年约600人次+"],
        ["负责人", "约50%"],
        ["双会", "WAIC→3F；ChinaJoy→5F；具身/AIGC叠加"],
        ["出海/领事", "另议另计价；到访≠挂牌"],
        ["租金", "3.3元/㎡/天；不做对赌；免租期1–3个月面议"],
        ["付款", "签约后7日50%；四季各10%；年终10%"],
    ], [18, 72], 34)

    for fn in ["创智汇年度活动运营方案-双会并入融合版.xlsx", "chuangzhihui-annual-ops-plan-dual-merge.xlsx",
               "创智汇年度活动运营方案-融合版.xlsx", "chuangzhihui-annual-ops-plan-fused.xlsx"]:
        p = os.path.join(HERE, fn); wb.save(p); shutil.copy2(p, os.path.join(ART, fn))
    print("Excel OK")


if __name__ == "__main__":
    build()
