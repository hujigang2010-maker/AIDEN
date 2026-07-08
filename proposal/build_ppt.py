# -*- coding: utf-8 -*-
"""创智汇 6600㎡ AI+IP 产业创新中心合作方案 · 高端商务版 PPT（完整版）。

包含：平面图（3F/5F）、杨浦政策导入与招商闭环、火山引擎园区政策、
三年合作条款与合作报价、媒体报价、规划建议与业务拓展等。
用法: python3 build_ppt.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
FONT = "Microsoft YaHei"
FONT_EN = "Arial"

# ===== 高端紫金主题（Royal Purple + Champagne Gold）=====
BG_A = "1A1140"; BG_B = "3A206E"          # 深紫渐变背景
PANEL = RGBColor(0x27, 0x1A, 0x4D); PANEL2 = RGBColor(0x1E, 0x15, 0x42)
INK = RGBColor(0xF5, 0xF1, 0xFF); MUT = RGBColor(0xC6, 0xBA, 0xEA); SOFT = RGBColor(0x94, 0x86, 0xC4)
ACC = RGBColor(0x8B, 0x7B, 0xFF); ACC2 = RGBColor(0xC7, 0x7D, 0xFF)   # 蓝紫 / 兰紫
GOLD = RGBColor(0xE9, 0xC2, 0x7C); GREEN = RGBColor(0x56, 0xD6, 0xB6); ROSE = RGBColor(0xF5, 0x8B, 0xB8)
LINE = RGBColor(0x4A, 0x3C, 0x7C)
H_ACC = "8B7BFF"; H_ACC2 = "C77DFF"; H_GOLD = "E9C27C"; H_GREEN = "56D6B6"; H_ROSE = "F58BB8"
DARKTX = RGBColor(0x24, 0x14, 0x40)       # 徽标上的深色字

SW, SH = Inches(13.333), Inches(7.5)
ML = Inches(0.85); CW = Inches(11.63)

prs = Presentation()
prs.slide_width = SW; prs.slide_height = SH
BLANK = prs.slide_layouts[6]
_FOOT_RUNS = []


def _no_shadow(sp): sp.shadow.inherit = False


def _set_gradient(shape, stops, angle_deg=90):
    spPr = shape._element.spPr
    for tag in ('a:noFill', 'a:solidFill', 'a:gradFill', 'a:blipFill', 'a:pattFill', 'a:grpFill'):
        e = spPr.find(qn(tag))
        if e is not None:
            spPr.remove(e)
    grad = spPr.makeelement(qn('a:gradFill'), {})
    gsLst = grad.makeelement(qn('a:gsLst'), {})
    for pos, color, alpha in stops:
        gs = grad.makeelement(qn('a:gs'), {'pos': str(int(pos * 1000))})
        clr = grad.makeelement(qn('a:srgbClr'), {'val': color})
        if alpha is not None:
            clr.append(grad.makeelement(qn('a:alpha'), {'val': str(int(alpha * 1000))}))
        gs.append(clr); gsLst.append(gs)
    grad.append(gsLst)
    grad.append(grad.makeelement(qn('a:lin'), {'ang': str(int(angle_deg * 60000)), 'scaled': '1'}))
    ln = spPr.find(qn('a:ln'))
    (ln.addprevious(grad) if ln is not None else spPr.append(grad))


def slide(gradient=True):
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    _no_shadow(bg); bg.line.fill.background()
    if gradient:
        _set_gradient(bg, [(0, BG_A, None), (55, BG_B, None), (100, BG_A, None)], 120)
    else:
        bg.fill.solid(); bg.fill.fore_color.rgb = DARKTX
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s


def rect(s, x, y, w, h, fill=None, line=None, lw=1.0, radius=False, grad=None, gang=90):
    t = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    b = s.shapes.add_shape(t, x, y, w, h); _no_shadow(b)
    if grad is not None:
        _set_gradient(b, grad, gang)
    elif fill is None:
        b.fill.background()
    else:
        b.fill.solid(); b.fill.fore_color.rgb = fill
    if line is None:
        b.line.fill.background()
    else:
        b.line.color.rgb = line; b.line.width = Pt(lw)
    return b


def oval(s, x, y, w, h, grad=None, gang=90, fill=None, line=None, lw=1.0):
    b = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h); _no_shadow(b)
    if grad is not None:
        _set_gradient(b, grad, gang)
    elif fill is not None:
        b.fill.solid(); b.fill.fore_color.rgb = fill
    else:
        b.fill.background()
    if line is None:
        b.line.fill.background()
    else:
        b.line.color.rgb = line; b.line.width = Pt(lw)
    return b


def text(s, x, y, w, h, runs, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, space=1.0, font=FONT, spacing=None):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    if isinstance(runs, str):
        runs = [(runs, color, bold)]
    p = tf.paragraphs[0]; p.alignment = align; p.line_spacing = space
    out = []
    for seg in runs:
        t, c, b = seg[0], seg[1], seg[2]
        f = seg[3] if len(seg) > 3 else font
        r = p.add_run(); r.text = t
        r.font.size = Pt(size); r.font.bold = b; r.font.color.rgb = c; r.font.name = f
        if spacing is not None:
            r._r.get_or_add_rPr().set('spc', str(int(spacing * 100)))
        out.append(r)
    return out


def bullets(s, x, y, w, h, items, size=14, color=MUT, gap=8, mark=ACC, lh=1.18):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(2); tf.margin_top = Pt(2)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = lh
        r0 = p.add_run(); r0.text = "▪  "; r0.font.size = Pt(size); r0.font.color.rgb = mark; r0.font.name = FONT
        r = p.add_run(); r.text = it; r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = FONT
    return tb


def header(s, sec, eyebrow, title):
    rect(s, ML, Inches(0.62), Inches(0.06), Inches(0.86), fill=GOLD)
    text(s, Emu(ML + Inches(0.22)), Inches(0.6), Inches(9), Inches(0.34), eyebrow, size=12.5, color=GOLD, bold=True, spacing=2, font=FONT_EN)
    text(s, Emu(ML + Inches(0.22)), Inches(0.92), Inches(10.4), Inches(0.66), title, size=26, color=INK, bold=True)
    text(s, Inches(11.3), Inches(0.5), Inches(1.55), Inches(0.9), [(sec, GOLD, True, FONT_EN)], size=40, align=PP_ALIGN.RIGHT, color=GOLD, font=FONT_EN)
    rect(s, Inches(11.55), Inches(1.32), Inches(1.28), Pt(2), fill=LINE)
    rect(s, ML, Inches(6.92), CW, Pt(1), fill=LINE)


def footer(s, idx):
    text(s, ML, Inches(7.0), Inches(9), Inches(0.32),
         [("上海创智汇 ", SOFT, False), ("CHUANGZHIHUI", SOFT, False, FONT_EN),
          ("  ·  AI+数字内容无界共创港 · 6600㎡ 谈判集中汇报", SOFT, False)], size=9.5)
    rs = text(s, Inches(11.2), Inches(7.0), Inches(1.63), Inches(0.32),
              [("%02d" % idx, GOLD, True, FONT_EN), (" / 00", SOFT, False, FONT_EN)],
              size=10.5, align=PP_ALIGN.RIGHT)
    _FOOT_RUNS.append(rs[1])


def card(s, x, y, w, h, title=None, body=None, items=None, accent=ACC, tcolor=INK,
         num=None, tsize=15.5, bsize=12.5):
    rect(s, x, y, w, h, grad=[(0, "2A1E55", None), (100, "1E1542", None)], gang=120, line=LINE, lw=1, radius=True)
    rect(s, x, y, Inches(0.07), h, fill=accent)
    ty = Emu(y + Inches(0.2)); tx = Emu(x + Inches(0.34))
    if num is not None:
        d = Inches(0.5)
        rect(s, tx, Emu(y + Inches(0.22)), d, d, grad=[(0, H_ACC, None), (100, H_ACC2, None)], gang=120, radius=True)
        text(s, tx, Emu(y + Inches(0.22)), d, d, num, size=16, color=DARKTX, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)
        tx = Emu(x + Inches(1.0))
    if title:
        text(s, tx, ty, Emu(w - (tx - x) - Inches(0.2)), Inches(0.5), title, size=tsize, color=tcolor, bold=True,
             anchor=MSO_ANCHOR.MIDDLE if num else MSO_ANCHOR.TOP)
    if body:
        text(s, Emu(x + Inches(0.34)), Emu(y + Inches(0.66)), Emu(w - Inches(0.6)), Emu(h - Inches(0.78)), body, size=bsize, color=MUT, space=1.16)
    if items:
        bullets(s, Emu(x + Inches(0.34)), Emu(y + Inches(0.62)), Emu(w - Inches(0.6)), Emu(h - Inches(0.72)), items, size=bsize, mark=accent)


def table(s, x, y, w, rows, col_w, sizes=None, rh=Inches(0.52), head_rh=Inches(0.5), first_col_color=GOLD):
    ncol = len(rows[0]); sizes = sizes or [12.5] * ncol; cy = y
    for ri, row in enumerate(rows):
        cur_h = head_rh if ri == 0 else rh
        if ri == 0:
            rect(s, x, cy, w, cur_h, grad=[(0, "36276C", None), (100, "2A1E55", None)], gang=0)
        elif ri % 2 == 0:
            rect(s, x, cy, w, cur_h, fill=RGBColor(0x22, 0x18, 0x44))
        cx = x
        for ci, cell in enumerate(row):
            cwid = Emu(int(w * col_w[ci]))
            tb = s.shapes.add_textbox(cx, cy, cwid, cur_h); tf = tb.text_frame; tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Pt(10); tf.margin_right = Pt(6); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
            p = tf.paragraphs[0]; p.line_spacing = 1.04
            r = p.add_run(); r.text = str(cell)
            if ri == 0:
                r.font.size = Pt(12.5); r.font.bold = True; r.font.color.rgb = GOLD
            else:
                r.font.size = Pt(sizes[ci]); r.font.bold = (ci == 0)
                r.font.color.rgb = first_col_color if ci == 0 else MUT
            r.font.name = FONT
            cx = Emu(cx + cwid)
        rect(s, x, Emu(cy + cur_h - Pt(0.75)), w, Pt(0.75), fill=LINE)
        cy = Emu(cy + cur_h)
    rect(s, x, y, Pt(2.5), Emu(cy - y), fill=GOLD)
    return cy


def zone(s, x, y, w, h, name, area, chex, sub=None):
    b = rect(s, x, y, w, h, radius=True, line=RGBColor.from_string(chex), lw=1.25)
    _set_gradient(b, [(0, chex, 24), (100, chex, 7)], 120)
    text(s, Emu(x + Inches(0.16)), Emu(y + Inches(0.12)), Emu(w - Inches(0.28)), Inches(0.5),
         name, size=12.5, color=INK, bold=True, space=1.05)
    text(s, Emu(x + Inches(0.16)), Emu(y + h - Inches(0.5)), Emu(w - Inches(0.28)), Inches(0.4),
         [(area, RGBColor.from_string(chex), True, FONT_EN)], size=15, font=FONT_EN)
    if sub:
        text(s, Emu(x + Inches(0.16)), Emu(y + Inches(0.44)), Emu(w - Inches(0.28)), Inches(0.4), sub, size=10, color=MUT)


def chev(s, x, y):
    text(s, x, y, Inches(0.5), Inches(0.6), [("▶", GOLD, True)], size=18, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def divider(idx, num, cn, en, desc):
    s = slide()
    oval(s, Inches(-2.2), Inches(2.4), Inches(7.5), Inches(7.5), grad=[(0, H_ACC, 16), (100, BG_A, 0)], gang=60)
    oval(s, Inches(9.2), Inches(-2.6), Inches(6.5), Inches(6.5), grad=[(0, H_ACC2, 16), (100, BG_A, 0)], gang=120)
    text(s, Inches(7.1), Inches(0.7), Inches(6.0), Inches(3.2), [(num, GOLD, True, FONT_EN)], size=200, align=PP_ALIGN.RIGHT, color=GOLD, font=FONT_EN)
    rect(s, ML, Inches(2.95), Inches(0.9), Pt(3), fill=GOLD)
    text(s, ML, Inches(3.2), Inches(9), Inches(0.4), [(en, GOLD, True, FONT_EN)], size=14, spacing=3)
    text(s, ML, Inches(3.6), Inches(10.5), Inches(1.1), cn, size=42, color=INK, bold=True)
    text(s, ML, Inches(4.85), Inches(10), Inches(0.5), desc, size=15, color=MUT, space=1.3)
    footer(s, idx)
    return s


# ============================================================
IDX = 0
def nxt():
    global IDX; IDX += 1; return IDX

# ---------- 1 封面 ----------
nxt(); s = slide()
oval(s, Inches(8.3), Inches(-2.0), Inches(7.5), Inches(7.5), grad=[(0, H_ACC, 22), (100, BG_A, 0)], gang=120)
oval(s, Inches(9.6), Inches(3.2), Inches(5.5), Inches(5.5), grad=[(0, H_ACC2, 18), (100, BG_A, 0)], gang=120)
rect(s, 0, 0, Inches(0.16), SH, grad=[(0, H_GOLD, None), (100, H_ACC2, None)], gang=90)
text(s, ML, Inches(0.95), Inches(11), Inches(0.4), [("SHANGHAI CHUANGZHIHUI · WUJIAOCHANG", GOLD, True, FONT_EN)], size=13, spacing=3)
rect(s, ML, Inches(1.52), Inches(0.9), Pt(2.5), fill=GOLD)
text(s, ML, Inches(1.78), Inches(11.6), Inches(1.7),
     [("上海创智汇", INK, True), ("\nAI + 数字内容无界共创港", GOLD, True)], size=40, bold=True, space=1.14)
text(s, ML, Inches(3.62), Inches(11.6), Inches(0.5),
     [("6600㎡ 招商运营合作", MUT, False), ("　·　", SOFT, False), ("谈判集中汇报", MUT, False)], size=16, spacing=1)
text(s, ML, Inches(4.2), Inches(11.6), Inches(0.9),
     [("五角场片区城市更新 · 创智汇（一期）OPC + AI + IP 创新中心\n", MUT, False),
      ("3 楼 孵化器 + 办公 ≈ ", MUT, False), ("2850㎡", GOLD, True),
      ("      5 楼 展厅 + 贸易 ≈ ", MUT, False), ("3670㎡", GOLD, True)], size=15, space=1.5)
tags = ["以 AI + 数字内容 为产业主轴", "政策导入 · 政企联动闭环", "三年合作 · 协同复制"]
tx = ML
for t in tags:
    wd = Inches(0.42 + len(t) * 0.14)
    rect(s, tx, Inches(5.5), wd, Inches(0.52), line=LINE, lw=1.25, radius=True)
    text(s, tx, Inches(5.5), wd, Inches(0.52), t, size=11.5, color=INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tx = Emu(tx + wd + Inches(0.2))
rect(s, ML, Inches(6.5), CW, Pt(1), fill=LINE)
text(s, ML, Inches(6.62), Inches(11), Inches(0.4), [("v3.0 谈判稿", GOLD, True), ("    项目定位（三选一）：AI+数字内容无界共创港 / AI+数据智能产业创新中心 / AI+在线新经济产业创新港", SOFT, False)], size=10.5)

# ---------- 2 议程 ----------
nxt(); s = slide(); header(s, "", "AGENDA", "汇报议程 · 九大板块"); footer(s, IDX)
ag = [("00", "研判与关键信息", "研判报告 · 事实底盘"), ("01", "项目与共识", "定位 / 双重身份 / 要点"),
      ("02", "空间与平面图", "3F 孵化办公 · 5F 展厅贸易"), ("03", "我们带来什么", "六大资源资产化"),
      ("04", "政策与闭环", "政策导入 · 产业落位 · 招商闭环"), ("05", "招商与运营", "漏斗 / 运营计划 / 活动"),
      ("06", "节点与排期", "三年 · 80/85/95% 去化"), ("07", "条款与报价", "三年 · 五类收费 · 测算"),
      ("08", "出海与国际会客厅", "国家会客厅 · 出海撮合"), ("09", "规划与拓展", "规划建议 · 协同复制"),
      ("10", "投决与下一步", "三问 · 30/60/90 天")]
cw, ch = Inches(2.82), Inches(1.5); gx, gy = Inches(0.13), Inches(0.18)
for i, (n, t, d) in enumerate(ag):
    x = Emu(ML + (i % 4) * (cw + gx)); y = Emu(Inches(1.9) + (i // 4) * (ch + gy))
    card(s, x, y, cw, ch, accent=(GOLD if i % 2 else ACC))
    text(s, Emu(x + Inches(0.3)), Emu(y + Inches(0.18)), Inches(1.5), Inches(0.55), [(n, GOLD if i % 2 else ACC, True, FONT_EN)], size=24, font=FONT_EN)
    text(s, Emu(x + Inches(0.3)), Emu(y + Inches(0.72)), Emu(cw - Inches(0.5)), Inches(0.4), t, size=13.5, color=INK, bold=True)
    text(s, Emu(x + Inches(0.3)), Emu(y + Inches(1.12)), Emu(cw - Inches(0.5)), Inches(0.4), d, size=10.5, color=MUT)

# ---------- 项目定位 ----------
nxt(); s = slide(); header(s, "", "POSITIONING", "项目定位 · 产业定位方向"); footer(s, IDX)
text(s, ML, Inches(1.72), Inches(11), Inches(0.35), "项目命名（三选一）", size=14, color=GOLD, bold=True)
names = [("方案一 · 推荐", "上海创智汇 · AI + 数字内容无界共创港", "最贴合「AI+数字内容」主轴与共创理念", GOLD, True),
         ("方案二", "上海创智汇 · AI + 数据智能产业创新中心", "偏 AI 数据智能与硬科技定位", ACC, False),
         ("方案三", "上海创智汇 · AI + 在线新经济产业创新港", "对齐杨浦「在线新经济」政策品牌", ACC2, False)]
for i, (tag, nm, d, c, rec) in enumerate(names):
    y = Emu(Inches(2.15) + i * Inches(0.92))
    rect(s, ML, y, CW, Inches(0.8), grad=[(0, "2A1E55", None), (100, "1E1542", None)], gang=120,
         line=(GOLD if rec else LINE), lw=(1.5 if rec else 1), radius=True)
    rect(s, ML, y, Inches(0.07), Inches(0.8), fill=c)
    text(s, Emu(ML + Inches(0.3)), y, Inches(2.1), Inches(0.8), [(tag, c, True)], size=13, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Emu(ML + Inches(2.5)), y, Inches(6.3), Inches(0.8), [(nm, INK, True)], size=17, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Emu(ML + Inches(8.7)), y, Inches(2.8), Inches(0.8), d, size=10.5, color=MUT, anchor=MSO_ANCHOR.MIDDLE, space=1.1)
rect(s, ML, Inches(5.05), CW, Inches(1.55), grad=[(0, "3A2878", None), (100, "2A1E55", None)], gang=0, line=GOLD, lw=1.25, radius=True)
text(s, Inches(1.1), Inches(5.18), Inches(11.2), Inches(0.4), [("产业定位方向　", GOLD, True), ("（经商讨暂定）", SOFT, False)], size=13.5)
text(s, Inches(1.1), Inches(5.6), Inches(11.2), Inches(0.95),
     [("以 ", MUT, False), ("人工智能和数字内容", GOLD, True), ("为产业主轴，重点承接 ", MUT, False),
      ("AI 微短剧、AIGC 内容生产、数字 IP、创意设计、智能营销、在线新经济、OPC 超级个体", INK, True),
      (" 等应用场景；并通过 办公空间 · 展示交易 · 订单撮合 · 政策服务 · 高校大厂资源 · IP 资源 · 出海服务，为入驻企业提供从落地到增长的综合服务。", MUT, False)],
     size=12.5, space=1.28)

# ---------- 3 关键信息 ----------
nxt(); s = slide(); header(s, "00", "FACT SHEET", "关键信息提取 · 项目基本盘"); footer(s, IDX)
y0 = Inches(1.78)
rect(s, ML, y0, Inches(6.05), Inches(3.05), grad=[(0, "2A1E55", None), (100, "1E1542", None)], gang=120, line=LINE, radius=True)
table(s, Emu(ML + Inches(0.18)), Emu(y0 + Inches(0.2)), Inches(5.7), [
    ["项目", "内容"], ["位置", "上海杨浦五角场 · 创智汇（创智天地）"], ["合作面积", "约 6600㎡"],
    ["3 楼 ≈2850㎡", "孵化器 + 办公（AI / OPC 主轴）"], ["5 楼 ≈3670㎡", "展厅 + 贸易（IP 内容主轴）"],
    ["合作周期", "3 年 · 可复制 · 协同周边 ≤2 项目"]], [0.30, 0.70], sizes=[12, 12], rh=Inches(0.46), head_rh=Inches(0.44))
card(s, Inches(7.15), y0, Inches(5.33), Inches(3.05), "区位与政策腹地", accent=GOLD, items=[
    "环同济：复旦 / 同济 / 财大 / 上理工高校群",
    "正处杨浦「YOUNG立方」内容集聚区核心圈",
    "紧邻 V聚场(6号楼) · 大学路 · B站新世代产业园",
    "离我方极近（20–30min）→ 轻驻场 · 高频次"])
yy = Inches(5.05)
card(s, ML, yy, Inches(3.72), Inches(1.72), "定位主线", accent=ACC,
     body="超级链接器：IP 为内容、AI 为工具、空间为载体；产业筋骨 · 文化灵魂 · 商业血脉。")
card(s, Inches(4.78), yy, Inches(3.72), Inches(1.72), "三大集群", accent=GREEN,
     body="动漫 IP  ·  科技应用  ·  交互设计（对齐杨浦「在线新经济·创意设计」千亿集群）", bsize=11.5)
card(s, Inches(8.76), yy, Inches(3.72), Inches(1.72), "可链接资源", accent=ACC2,
     body="北大上海校友会、同济设计创新院、科企联、IP/玩具/广告协会、混知等 IP、聚成智能、中建四局。", bsize=11)

# ---------- 项目研判报告 · 市场与区位 ----------
nxt(); s = slide(); header(s, "00", "RESEARCH", "项目研判报告 · 市场与区位"); footer(s, IDX)
rr = [("区位研判", GOLD, "五角场城市副中心 + 环同济经济圈；创智天地正处杨浦 YOUNG立方内容集聚区核心圈，紧邻 V聚场/大学路/B站新世代产业园。"),
      ("市场研判", ACC, "AI+ 与数字内容处政策红利期：微短剧/AIGC/数字人爆发；上海明确徐汇·杨浦·闵行建 AI 微短剧集聚区，需求侧旺盛。"),
      ("政策研判", ACC2, "市区两级政策叠加（YOUNG立方/长阳秀带/人工智能+/微短剧/模塑申城/火山引擎）；载体认定 + 三券 + 补贴可形成兑现能力。"),
      ("客群研判", GREEN, "AI 应用/AIGC/数字人、AI 短剧、数字 IP、创意设计、智能营销、OPC 超级个体；供应链（玩具/毛绒/潮玩）与内容 IP 出海需求明确。")]
for i, (t, c, d) in enumerate(rr):
    x = Emu(ML + (i % 2) * Inches(5.92)); y = Emu(Inches(1.85) + (i // 2) * Inches(1.72))
    card(s, x, y, Inches(5.72), Inches(1.55), title=t, body=d, accent=c, bsize=12)
rect(s, ML, Inches(5.35), CW, Inches(1.3), grad=[(0, "3A2878", None), (100, "2A1E55", None)], gang=0, line=GOLD, lw=1, radius=True)
text(s, Inches(1.1), Inches(5.48), Inches(11.2), Inches(0.4), [("研判结论　", GOLD, True), ("需求真实、政策强、区位优——项目具备「政策引导产业落位 + 出海撮合放大」的稀缺条件。", INK, False)], size=13)
bullets(s, Inches(1.1), Inches(5.92), Inches(11), Inches(0.7),
        ["6600㎡ 体量适中，可做「小而精」样板并向周边复制（≤2 项目）",
         "对标魔术空间等：高频活动（约 110 场级）驱动招商与载体 KPI"], size=11, gap=3, mark=GOLD)

# ---------- 项目研判报告 · SWOT ----------
nxt(); s = slide(); header(s, "00", "RESEARCH · SWOT", "项目研判报告 · SWOT 与对策"); footer(s, IDX)
sw = [("S 优势", GREEN, ["政策核心圈 + 环同济 + 高校大厂资源", "AI/OPC 社群 + IP 内容 + 领事/出海资源", "离运营方近，轻驻场高频次"]),
      ("W 劣势", ROSE, ["体量小（6600㎡），单靠租金想象空间有限", "载体/OPC 社区等「经认定」资格待落实", "出海撮合需自有资源与专职团队"]),
      ("O 机会", ACC, ["微短剧/AIGC 政策红利与集聚区身份", "出海撮合 + 国家会客厅高毛利业务", "沪穗联动、可复制到周边项目"]),
      ("T 威胁", ACC2, ["招商去化对赌压力（三年均 85%）", "政策以最新文件为准，兑现有不确定性", "出海行业分散、跨行业撮合难"])]
for i, (t, c, items) in enumerate(sw):
    x = Emu(ML + (i % 2) * Inches(5.92)); y = Emu(Inches(1.82) + (i // 2) * Inches(2.32))
    card(s, x, y, Inches(5.72), Inches(2.15), title=t, items=items, accent=c, bsize=11.5)

# ---------- 4 章节01 ----------
nxt(); divider(IDX, "01", "项目与共识", "PROJECT & CONSENSUS", "不做二房东，做产业空间合作运营")

# ---------- 5 共识 ----------
nxt(); s = slide(); header(s, "01", "PROJECT & CONSENSUS", "双重身份 + 三件谈判要点"); footer(s, IDX)
card(s, ML, Inches(1.8), Inches(5.72), Inches(1.5), "身份一 · 产业孵化核心", accent=ACC,
     body="五角场城市更新的孵化核心与先行示范区——空间改造 + 产业升级 + 商业激活的完整落地路径。")
card(s, Inches(6.76), Inches(1.8), Inches(5.72), Inches(1.5), "身份二 · 超级链接器", accent=ACC2,
     body="环同济经济圈资源链接平台：IP 内容 × AI 工具 × 空间载体，闭环生长。")
text(s, ML, Inches(3.55), Inches(8), Inches(0.4), "本次谈判要点", size=17, color=GOLD, bold=True)
pts = [("1", "合作模式", "对赌运营 · 收基础运营费\n+ 招商佣金 + 增值分成"),
       ("2", "商业条款", "3 年合作 · 五类收费\n协同周边 ≤2 项目 · 可复制"),
       ("3", "启动时间", "与装修 / 设计交付同步，\n尽快锁定并启动招商")]
for i, (n, t, d) in enumerate(pts):
    x = Emu(ML + i * Inches(3.94))
    card(s, x, Inches(4.05), Inches(3.72), Inches(1.5), num=n, title=t, accent=GOLD)
    text(s, Emu(x + Inches(0.34)), Emu(Inches(4.05) + Inches(0.8)), Inches(3.3), Inches(0.66), d, size=11.5, color=MUT, space=1.18)
rect(s, ML, Inches(5.78), CW, Inches(1.0), grad=[(0, "322462", None), (100, "221848", None)], gang=0, line=GOLD, lw=1, radius=True)
text(s, Inches(1.1), Inches(5.78), Inches(11.1), Inches(1.0),
     [("核心共识　", GOLD, True), ("把 6600㎡ 升级为「产业空间合作运营项目」——我方输出 政策 + 招商 + 产业资源 + 活动 + 媒体 五项能力，按「基础运营费 + 招商佣金 + 增值分成」取酬。", INK, False)],
     size=13.5, anchor=MSO_ANCHOR.MIDDLE, space=1.2)

# ---------- 6 章节02 空间 ----------
nxt(); divider(IDX, "02", "空间与平面图", "SPACE & FLOOR PLAN", "3F 孵化办公（AI 主轴）· 5F 展厅贸易（IP 主轴）")

# ---------- 7 3F 平面图（按真实布局复原） ----------
nxt(); s = slide(); header(s, "02", "FLOOR PLAN · 3F", "3 楼平面布局 · OPC + AI + IP 创新中心 ≈ 2850㎡"); footer(s, IDX)
fx, fy, fw, fh = ML, Inches(1.78), Inches(8.05), Inches(4.85)
rect(s, fx, fy, fw, fh, line=GOLD, lw=1.5, radius=True, grad=[(0, "231A4A", None), (100, "1A1236", None)], gang=120)
pad = Inches(0.18)
ix, iy = Emu(fx + pad), Emu(fy + pad); iw = Emu(fw - pad * 2); ih = Emu(fh - pad * 2)
lw = Emu(iw * 0.56); gp = Inches(0.16); rx = Emu(ix + lw + gp); rwd = Emu(iw - lw - gp)
# 左列：AI+IP 600 / 配套 / AI+IP 1150
zone(s, ix, iy, lw, Inches(1.28), "AI+IP 产业办公空间", "600㎡", H_GOLD, sub="316 (150) · 301 (175)")
zone(s, ix, Emu(iy + Inches(1.44)), lw, Inches(0.72), "接待 · 前台 · 休闲区 · 茶水间", "配套", H_ROSE)
zone(s, ix, Emu(iy + Inches(2.32)), lw, Inches(2.02), "AI+IP 产业办公空间", "1150㎡", H_GOLD, sub="306·305·304·303·302·310")
# 右列：直播间 250 / AI展示 150 / OPC 700
zone(s, rx, iy, rwd, Inches(1.02), "直播间", "250㎡", H_ACC)
zone(s, rx, Emu(iy + Inches(1.18)), rwd, Inches(0.86), "AI 展示与运营", "150㎡", H_GREEN)
zone(s, rx, Emu(iy + Inches(2.2)), rwd, Inches(2.14), "OPC 产业办公空间", "700㎡", H_ACC2)
card(s, Inches(9.15), Inches(1.78), Inches(3.33), Inches(4.85), "楼层要点", accent=GOLD, items=[
    "OPC + AI + IP 创新中心整层运营",
    "9 个可分割办公单元（301–316）",
    "含 AI 展示与运营 + 直播间双运营位",
    "主轴：AI / OPC — 对齐「人工智能+」",
    "适配：AI 初创 / 研发设计 / 孵化项目",
    "设计建议：4–6 人设计条件，不做联合办公"])
text(s, ML, Inches(6.68), Inches(8), Inches(0.3), "＊依据一期方案平面图复原（分区/房号/面积对应）；最终以物业实测与设计院图纸为准。", size=9.5, color=SOFT)

# ---------- 8 5F 平面图（按真实布局复原） ----------
nxt(); s = slide(); header(s, "02", "FLOOR PLAN · 5F", "5 楼平面布局 · 综合集群展厅 ≈ 3670㎡"); footer(s, IDX)
fx, fy, fw, fh = ML, Inches(1.78), Inches(8.05), Inches(4.85)
rect(s, fx, fy, fw, fh, line=GOLD, lw=1.5, radius=True, grad=[(0, "231A4A", None), (100, "1A1236", None)], gang=120)
ix, iy = Emu(fx + pad), Emu(fy + pad); iw = Emu(fw - pad * 2); ih = Emu(fh - pad * 2)
g = Inches(0.12)
cA = Emu(iw * 0.205); cB = Emu(iw * 0.09); cC = Emu(iw * 0.225); cD = Emu(iw * 0.225); cE = Emu(iw * 0.205)
xA = ix; xB = Emu(xA + cA + g); xC = Emu(xB + cB + g); xD = Emu(xC + cC + g); xE = Emu(xD + cD + g)
zone(s, xA, iy, cA, ih, "AI+IP 展示中心", "700㎡", H_ROSE, sub="景区+博物馆文创 · 50+ IP 展示")
zone(s, xB, iy, cB, ih, "休闲沙龙区", "", H_ACC)
zone(s, xC, iy, cC, Inches(1.3), "东莞潮玩产业集群", "330㎡", H_GREEN)
zone(s, xC, Emu(iy + Inches(1.42)), cC, Inches(3.02), "汕头玩具产业集群", "900㎡", H_ACC2)
zone(s, xD, iy, cD, Inches(1.12), "AI 智能产业集群", "270㎡", H_GOLD)
zone(s, xD, Emu(iy + Inches(1.24)), cD, Inches(3.2), "扬州毛绒产业集群", "700㎡", H_GOLD)
zone(s, xE, iy, cE, Inches(0.8), "仓储", "120㎡", "9486C4")
zone(s, xE, Emu(iy + Inches(0.92)), cE, Inches(3.52), "OPC 社区培训中心", "650㎡", H_ACC)
card(s, Inches(9.15), Inches(1.78), Inches(3.33), Inches(4.85), "楼层要点", accent=GOLD, items=[
    "50+ IP 展示 + 跨境展销贸易",
    "四大集群：汕头玩具·扬州毛绒·东莞潮玩·文创",
    "顶层可办漫展 / 展览（门票+赞助）",
    "主轴：IP 内容 — 对齐「互联网内容创作」",
    "适配：IP 品牌方 / 玩具外贸 / 文创贸易",
    "培训 + 沙龙 + 小红书招募 → 经营性收入"])
text(s, ML, Inches(6.68), Inches(8), Inches(0.3), "＊依据一期方案平面图复原（分区/面积对应）；最终以物业实测与设计院图纸为准。", size=9.5, color=SOFT)

# ---------- 9 章节03 ----------
nxt(); divider(IDX, "03", "我们能带来什么", "WHAT WE BRING", "六大稀缺资源，全部可资产化为招商工具")

# ---------- 10 六大资源 ----------
nxt(); s = slide(); header(s, "03", "WHAT WE BRING", "资源即招商 · 六大引擎"); footer(s, IDX)
res = [("① 政策抓手", "杨浦 YOUNG立方 / 长阳秀带 / 人工智能+ + 火山引擎园区政策 + 算力补贴 →「创智汇专有政策包」", GOLD),
       ("② 高校 / 协会背书", "北大上海校友会、同济设计创新院、复旦、科企联、IP/玩具/广告协会 → 挂牌即招商", ACC),
       ("③ IP 内容资源", "混知、观山礼玉等 IP 品牌方 + 汕头玩具/扬州毛绒/东莞潮玩/文创集群 → 5 楼自带内容", ACC2),
       ("④ AI / OPC 社群", "自有 OPC（AI 开放社群）品牌、黑客松、AI 项目库、聚成智能 → 3 楼自带客流", GREEN),
       ("⑤ 活动 / 会客厅", "漫展（门票+赞助）、北欧会客厅（出海撮合）、沙龙、峰会借势 → 活动即招商", ACC),
       ("⑥ 资本 / 产业方", "江西金控基金、复旦科技园导流、中建四局工程资源、奇瑞/华为外延 → 链主带动", ROSE)]
cw2 = Inches(3.78); ch2 = Inches(1.72)
for i, (t, d, c) in enumerate(res):
    x = Emu(ML + (i % 3) * (cw2 + Inches(0.14))); y = Emu(Inches(1.8) + (i // 3) * (ch2 + Inches(0.2)))
    card(s, x, y, cw2, ch2, title=t, body=d, accent=c, bsize=11.5)
rect(s, ML, Inches(5.95), CW, Inches(0.72), grad=[(0, "322462", None), (100, "221848", None)], gang=0, line=LINE, radius=True)
text(s, Inches(1.1), Inches(5.95), Inches(11), Inches(0.72),
     [("普通园区「发广告 · 等上门」；六维齐发 → 客户 ", MUT, False), ("「被送进来」而非「被拉进来」。", GOLD, True)], size=13.5, anchor=MSO_ANCHOR.MIDDLE)

# ---------- 11 章节04 政策与闭环 ----------
nxt(); divider(IDX, "04", "政策导入与招商闭环", "POLICY & CLOSED LOOP", "以政策引导产业布局落位，形成产业招商闭环")

# ---------- 12 杨浦政策矩阵 ----------
nxt(); s = slide(); header(s, "04", "YANGPU POLICY", "杨浦政策矩阵 · 项目正处政策核心圈"); footer(s, IDX)
pol = [("YOUNG立方 · 18 条", "互联网优质内容 / 微短剧",
        "免：入孵 V聚场 三年租金物业减免（全免/75%/50%）；补：直播券≤50万、活动≤50%（年≤200万）；投：≥10亿内容基金、贷款50%贴息；人才：购房≤200万、公寓≤8000/月", GOLD),
       ("长阳秀带 · 在线新经济", "数字经济 / AI 大数据",
        "AI 及大数据企业房租补贴 3 年每年≤100万；头部/总部≤1500万开办费；技术攻关研发 10–30%≤200万；应用场景≤500万；10 亿母基金", ACC),
       ("人工智能+ 行动（26–28）", "六大赛道 · 内容生成",
        "大模型垂类应用补贴≤50%、≤500万；AI 人才个人奖励≤30万；聚焦具身智能/智能终端/内容生成工具等", ACC2),
       ("数字经济产业体系", "三大千亿集群",
        "在线新经济 · 智能制造 · 创意设计 三大千亿集群 + 四大新兴；专项资金「就高不重复」；算力补贴叠加", GREEN)]
for i, (t, tag, d, c) in enumerate(pol):
    x = Emu(ML + (i % 2) * Inches(5.92)); y = Emu(Inches(1.8) + (i // 2) * Inches(1.72))
    card(s, x, y, Inches(5.72), Inches(1.55), accent=c)
    text(s, Emu(x + Inches(0.34)), Emu(y + Inches(0.16)), Inches(3.9), Inches(0.4), t, size=14.5, color=INK, bold=True)
    text(s, Emu(x + Inches(0.34)), Emu(y + Inches(0.52)), Inches(5.1), Inches(0.3), [(tag, c, True)], size=10.5)
    text(s, Emu(x + Inches(0.34)), Emu(y + Inches(0.82)), Inches(5.15), Inches(0.66), d, size=10.3, color=MUT, space=1.14)
rect(s, ML, Inches(5.35), CW, Inches(1.32), grad=[(0, "3A2878", None), (100, "2A1E55", None)], gang=0, line=GOLD, lw=1, radius=True)
text(s, Inches(1.1), Inches(5.5), Inches(11.2), Inches(0.4), [("关键卡位　", GOLD, True), ("创智汇位于创智天地——正是 YOUNG立方「一楼(V聚场6号楼) · 一街(大学路) · 一园(B站新世代产业园) · 一区」核心腹地。", INK, False)], size=12.5)
text(s, Inches(1.1), Inches(5.95), Inches(11.2), Inches(0.6), [("　　　　　　", GOLD, True), ("政策适配度极高：3F AI 孵化对齐「人工智能+」，5F IP 内容 / 微短剧对齐「互联网内容创作」，可直接承接补贴与载体认定。", MUT, False)], size=11.5, space=1.15)
text(s, ML, Inches(6.72), Inches(11.6), Inches(0.2), "＊政策口径以杨浦区政府最新发布文件为准，本页为要点提炼。", size=9, color=SOFT)

# ---------- 13 招商闭环 ----------
nxt(); s = slide(); header(s, "04", "CLOSED LOOP", "政策 → 产业布局落位 → 招商闭环"); footer(s, IDX)
steps = [("01", "政策导入", "YOUNG立方 / 长阳秀带 /\n人工智能+ / 火山引擎", GOLD),
         ("02", "产业布局落位", "3F AI 孵化 · 5F IP 展贸\n按政策赛道分区落位", ACC),
         ("03", "精准招商", "牌照锚定 · 政策礼包 ·\n社群/活动/复旦科技园导流", ACC2),
         ("04", "企业成长服务", "联合公司：注册 / 申报 /\n知产 / 融资 / 出海撮合", GREEN),
         ("05", "补贴资质反哺", "补贴归运营平台 · 载体认定\n→ 提升招商可信度", ROSE)]
cwf = Inches(2.1); y = Inches(2.05)
for i, (n, t, d, c) in enumerate(steps):
    x = Emu(ML + i * (cwf + Inches(0.22)))
    card(s, x, y, cwf, Inches(2.55), accent=c)
    text(s, Emu(x + Inches(0.26)), Emu(y + Inches(0.22)), Inches(1.2), Inches(0.6), [(n, c, True, FONT_EN)], size=24, font=FONT_EN)
    text(s, Emu(x + Inches(0.26)), Emu(y + Inches(0.85)), Emu(cwf - Inches(0.4)), Inches(0.5), t, size=13.5, color=INK, bold=True)
    text(s, Emu(x + Inches(0.26)), Emu(y + Inches(1.3)), Emu(cwf - Inches(0.4)), Inches(1.1), d, size=10, color=MUT, space=1.15)
    if i < 4:
        chev(s, Emu(x + cwf), Emu(y + Inches(0.95)))
rect(s, ML, Inches(5.0), CW, Inches(1.55), grad=[(0, "322462", None), (100, "221848", None)], gang=0, line=GOLD, lw=1, radius=True)
text(s, Inches(1.1), Inches(5.15), Inches(11.2), Inches(0.4), [("闭环逻辑　", GOLD, True), ("政策吸引企业 → 企业落位对齐产业 → 高效招商去化 → 联合公司服务企业成长 → 补贴/资质反哺平台 → 再吸引更优企业。", INK, False)], size=12.5)
bullets(s, Inches(1.1), Inches(5.65), Inches(11), Inches(0.9),
        ["越招越准：政策把符合赛道的企业「筛」进来，去化质量高、租金承载强",
         "越做越强：补贴与载体资质回流平台，招商可信度与议价力持续提升"], size=11.5, gap=4, mark=GOLD)

# ---------- 14 火山引擎政策 ----------
nxt(); s = slide(); header(s, "04", "VOLCANO ENGINE", "火山引擎园区独立政策（优于火山工坊）"); footer(s, IDX)
card(s, ML, Inches(1.8), Inches(4.05), Inches(2.05), "一、无门槛半年免费", accent=GOLD,
     body="入园企业半年费用免费，按预估半年费用一次性发放代金券，无消费门槛。")
card(s, ML, Inches(4.0), Inches(4.05), Inches(2.5), "作为招商抓手", accent=GREEN, items=[
    "叠加杨浦算力补贴（凭票最高 50%）",
    "叠加腾讯云免 2 月 / 85 折",
    "打造「创智汇专有算力政策包」",
    "签约共建，作为 AI 企业入驻礼包"])
rect(s, Inches(5.1), Inches(1.8), Inches(7.38), Inches(0.55), grad=[(0, "36276C", None), (100, "2A1E55", None)], gang=0, radius=True)
text(s, Inches(5.3), Inches(1.8), Inches(7), Inches(0.55), [("二、大客户额外折扣（代金券之外，累计消费越多折扣越低）", GOLD, True)], size=13, anchor=MSO_ANCHOR.MIDDLE)
table(s, Inches(5.1), Inches(2.45), Inches(7.38), [
    ["累计消费", "额外折扣"], ["0 – 10 万", "5 折 ~ 7 折"], ["10 – 30 万", "4.5 折 ~ 5 折"],
    ["30 – 50 万", "4 折 ~ 4.5 折"], ["50 – 100 万", "3.5 折 ~ 4 折"], ["100 – 300 万", "3 折 ~ 3.5 折"],
    ["300 – 500 万+", "2.5 折 ~ 3 折"]], [0.5, 0.5], sizes=[13, 13], rh=Inches(0.5), head_rh=Inches(0.44), first_col_color=INK)

# ---------- 政策工具箱 · 三端赋能 ----------
nxt(); s = slide(); header(s, "04", "POLICY TOOLBOX", "人工智能+数字内容 · 政策工具箱（三端）"); footer(s, IDX)
tri = [("企业端 · 补给企业", GOLD, [
    "科技创新券：企业≤30万/年、团队≤10万/年（≤50%）",
    "模塑申城：算力券 / 模型券 / 语料券",
    "高企首认 20万；专精特新市级 10万 / 国家小巨人 30万",
    "AI 大数据房租 2.0 元/天/㎡（人均 15㎡）≤100万",
    "经营性奖励 + 参保奖励（10–20 人给 10万，超 20 每人≤1.5万）"]),
    ("载体端 · 强平台", ACC, [
    "科技服务业：成果转化服务平台≤100万/年运营补贴",
    "服务业引导资金：市级≥500万投资 / 区级 500–1000万，≤300万",
    "OPC 超级个体社区：AI 工具链 / 智能体 / Tokens券 / 大厂团队 各≤1000–2000万",
    "未来产业试验场 / 加速器（认定后申报）",
    "→ 把创智汇做成「经认定载体 + 成果转化服务平台」"]),
    ("场景端 · 造场景", ACC2, [
    "AI 微短剧集聚区（徐汇/杨浦/闵行）：智能体≤研发20%≤1000万",
    "优秀微短剧≤300万；优质厂牌≤200万；出海≤50万；课程≤20万",
    "标志性重大应用场景 / 创新产品市场化应用（联合≤5家采购）",
    "消费新场景专项：补助≤30%、≥100万、累计≤2000万",
    "→ 5楼做「AI 数字内容高价值应用场景」"])]
for i, (t, c, items) in enumerate(tri):
    x = Emu(ML + i * Inches(3.94))
    card(s, x, Inches(1.8), Inches(3.72), Inches(4.55), title=t, items=items, accent=c, bsize=10.8)
text(s, ML, Inches(6.5), Inches(11.6), Inches(0.3),
     "＊政策口径以各级政府最新发布文件为准；「经认定」类政策需先取得认定资格方可对外兑现。", size=9.5, color=SOFT)

# ---------- 政企联动 · 辅助招商定位 ----------
nxt(); s = slide(); header(s, "04", "GOV-ENTERPRISE SYNERGY", "政企联动 · 辅助招商定位"); footer(s, IDX)
rect(s, ML, Inches(1.75), CW, Inches(0.95), grad=[(0, "3A2878", None), (100, "2A1E55", None)], gang=0, line=GOLD, lw=1.25, radius=True)
text(s, Inches(1.1), Inches(1.75), Inches(11.2), Inches(0.95),
     [("招商定位　", GOLD, True), ("创智汇 = 杨浦「AI + 数字内容」孵化转化载体 + 场景清单发布平台（3F 孵化转化 · 5F 场景应用与交易）", INK, True)],
     size=15, anchor=MSO_ANCHOR.MIDDLE)
syn = [("企业端 → 招什么", GOLD, "招「能拿补贴的优质企业」：AI/AIGC、数字人、智能体、AI 微短剧、内容生成工具、科技型中小；用三券与房租/经营奖励做入驻礼包。"),
       ("载体端 → 靠什么", ACC, "把创智汇做成经认定载体 / OPC 社区 / 成果转化服务平台 / 服务业园区，取得政策兑现资格与运营补贴。"),
       ("场景端 → 卖什么", ACC2, "5楼做 AI 数字内容高价值应用场景、AI 微短剧承接点、创新产品市场化应用撮合，形成可背书的场景清单。")]
for i, (t, c, d) in enumerate(syn):
    x = Emu(ML + i * Inches(3.94))
    card(s, x, Inches(2.9), Inches(3.72), Inches(1.95), title=t, body=d, accent=c, bsize=11.5)
rect(s, ML, Inches(5.05), CW, Inches(1.5), grad=[(0, "2A1E55", None), (100, "1E1542", None)], gang=120, line=LINE, radius=True)
rect(s, ML, Inches(5.05), Inches(0.07), Inches(1.5), fill=GREEN)
text(s, Inches(1.1), Inches(5.18), Inches(11), Inches(0.4), "政企联动机制", size=14.5, color=GREEN, bold=True)
bullets(s, Inches(1.1), Inches(5.62), Inches(5.6), Inches(0.9),
        ["挂牌「杨浦科技企业服务中心」，接入政务快速通道",
         "对接区科经委 / 发改委 / 文旅局 / 市经信委"], size=11.5, gap=4, mark=GREEN)
bullets(s, Inches(6.9), Inches(5.62), Inches(5.5), Inches(0.9),
        ["建企业「三券」适配表，联合公司承接申报",
         "补贴归运营平台公司，反哺招商可信度"], size=11.5, gap=4, mark=GREEN)

# ---------- 后续摸排清单 ----------
nxt(); s = slide(); header(s, "04", "NEXT-STEP CHECKLIST", "后续摸排清单 · 政企联动落地"); footer(s, IDX)
table(s, ML, Inches(1.8), CW, [
    ["摸排事项", "对接部门", "目的 / 判定"],
    ["能否按创新创业载体 / 成果转化服务平台 / 场景清单申报", "区科经委", "确立载体端政策兑现资格"],
    ["3F 办公孵化 + 5F 展示交易能否作服务业引导资金项目", "区发改委", "核算投资是否达 500 万门槛"],
    ["取得 OPC 超级个体社区「认定条件」", "区人民政府 / 科经委", "先决条件；未认定只能做协同点"],
    ["杨浦 AI 微短剧集聚区承接空间；5F 可否作展示/路演/拍摄/协同点", "区文旅局", "确立场景端承接身份"],
    ["建企业「算力/模型/语料」三券适配表", "市经信委 / 平台", "核查合同 / 账单 / 发票 / 使用日志"],
    ["科企联 / 服务机构是否入驻创新券平台", "科企联 / 服务机构", "打通创新券在线下单核验兑付"],
    ["入驻企业产品进《创新产品推荐目录》→ 5F 找采购/承租联合申报", "市经信委 / 采购方", "落地 5F 订单撮合机制"]],
    [0.5, 0.2, 0.3], sizes=[11.5, 11, 11], rh=Inches(0.56), head_rh=Inches(0.46), first_col_color=GOLD)

# ---------- 政企联动 · 生态闭环（双环） ----------
nxt(); s = slide(); header(s, "04", "ECO CLOSED LOOP", "政企联动 · 生态闭环（双环驱动）"); footer(s, IDX)
outer = [("政府政策", "YOUNG立方/长阳秀带/人工智能+/微短剧/火山引擎", GOLD),
         ("平台承接", "载体认定 · 挂牌 · 成果转化服务平台", ACC),
         ("精准招商", "政策礼包 · 三券适配 · 场景背书", ACC2),
         ("企业成长", "联合公司：注册/申报/知产/融资/出海", GREEN),
         ("场景与内容", "5F 微短剧/展贸/IP · 创新产品应用", ROSE),
         ("补贴反哺", "补贴归平台 · 资质回流 · 能级提升", GOLD)]
for i, (t, d, c) in enumerate(outer):
    x = Emu(ML + (i % 3) * Inches(3.94)); y = Emu(Inches(1.8) + (i // 3) * Inches(1.55))
    card(s, x, y, Inches(3.72), Inches(1.4), title=t, body=d, accent=c, tsize=14.5, bsize=11)
    if i % 3 != 2:
        chev(s, Emu(x + Inches(3.72)), Emu(y + Inches(0.45)))
rect(s, ML, Inches(4.95), CW, Inches(1.6), grad=[(0, "322462", None), (100, "221848", None)], gang=0, line=GOLD, lw=1, radius=True)
text(s, Inches(1.1), Inches(5.08), Inches(11.2), Inches(0.4), [("双环驱动　", GOLD, True), ("① 政企联动闭环：政府—平台—企业—场景—补贴反哺，能级持续提升；② 产业生态闭环：内容引流—产业承接—商业变现—服务成长。", INK, False)], size=12)
bullets(s, Inches(1.1), Inches(5.55), Inches(11), Inches(0.9),
        ["政企联动：把政策转化为招商资格与兑现能力，越做越有背书",
         "生态自进化：内容 × 产业 × 商业 × 服务相互反哺，摆脱单一租金依赖"], size=11.5, gap=4, mark=GOLD)

# ---------- 15 章节05 招商 ----------
nxt(); divider(IDX, "05", "招商方案", "LEASING STRATEGY", "IP + AI 双轨 · 楼层产品化 · 四级漏斗")

# ---------- 16 楼层定位 ----------
nxt(); s = slide(); header(s, "05", "LEASING STRATEGY", "楼层产品切分 · 招商最小颗粒度"); footer(s, IDX)
rect(s, ML, Inches(1.78), Inches(5.72), Inches(4.85), grad=[(0, "2A1E55", None), (100, "1E1542", None)], gang=120, line=LINE, radius=True)
rect(s, ML, Inches(1.78), Inches(0.07), Inches(4.85), fill=ACC)
text(s, Inches(1.1), Inches(1.95), Inches(5.4), Inches(0.4), [("3 楼", ACC, True), ("　约 2850㎡　孵化 + 办公", INK, True)], size=16)
table(s, Inches(1.05), Inches(2.5), Inches(5.4), [
    ["产品", "面积", "客户"], ["标准小单元", "80–150㎡×8–12", "AI 初创/设计/小微"],
    ["成长型单元", "200–350㎡×3–4", "AI 应用/研发服务"], ["OPC 联合办公", "工位大区", "AI 项目/黑客松/孵化"],
    ["直播 / AI 展示", "共享 1–2 间", "按次 / 时段"]],
    [0.33, 0.34, 0.33], sizes=[11.5, 11, 11], rh=Inches(0.46), head_rh=Inches(0.44), first_col_color=INK)
rect(s, Inches(6.76), Inches(1.78), Inches(5.72), Inches(4.85), grad=[(0, "2A1E55", None), (100, "1E1542", None)], gang=120, line=LINE, radius=True)
rect(s, Inches(6.76), Inches(1.78), Inches(0.07), Inches(4.85), fill=GOLD)
text(s, Inches(7.02), Inches(1.95), Inches(5.4), Inches(0.4), [("5 楼", GOLD, True), ("　约 3670㎡　展厅 + 贸易", INK, True)], size=16)
table(s, Inches(6.97), Inches(2.5), Inches(5.4), [
    ["产品", "面积", "内容"], ["综合集群展厅", "700–900㎡", "50+ IP 轮展"],
    ["产业集群展位", "270–700㎡×4–5", "玩具/毛绒/潮玩/文创"], ["IP 展销/跨境", "按摊位", "IP 零售/跨境展销"],
    ["培训沙龙/仓储", "120–650㎡", "OPC 培训/活动"]],
    [0.33, 0.34, 0.33], sizes=[11.5, 11, 11], rh=Inches(0.46), head_rh=Inches(0.44), first_col_color=INK)

# ---------- 17 漏斗 + 目标 ----------
nxt(); s = slide(); header(s, "05", "LEASING STRATEGY", "四级招商漏斗 + 三年去化目标"); footer(s, IDX)
fn = [("L1", "牌照锚定", "校友会/科企联/IP/玩具协会挂牌前置"),
      ("L2", "政策招商", "杨浦政策 + 火山引擎 + 算力补贴礼包"),
      ("L3", "社群带流", "OPC 社群/黑客松/小红书/漫展/复旦科技园"),
      ("L4", "资源转化", "高校成果转化/IP 品牌方/外贸玩具厂转化")]
cwf = Inches(2.78)
for i, (n, t, d) in enumerate(fn):
    x = Emu(ML + i * (cwf + Inches(0.16)))
    card(s, x, Inches(1.85), cwf, Inches(2.0), accent=ACC)
    text(s, Emu(x + Inches(0.34)), Inches(2.05), Inches(1.6), Inches(0.5), [(n, ACC, True, FONT_EN)], size=26, font=FONT_EN)
    text(s, Emu(x + Inches(0.34)), Inches(2.62), Inches(2.3), Inches(0.4), t, size=15, color=INK, bold=True)
    text(s, Emu(x + Inches(0.34)), Inches(3.08), Inches(2.4), Inches(0.7), d, size=10.5, color=MUT, space=1.12)
kpi = [("三年均 85%", "或 80% → 85% → 95%（首/次/三年）", GOLD),
       ("对赌挂钩", "基础运营费与去化目标对赌", ACC),
       ("5 类客群", "AI / 潮玩玩具文创 IP / 科技中小 /\n高校转化 / 生产性服务", GREEN)]
for i, (v, l, c) in enumerate(kpi):
    x = Emu(ML + i * Inches(3.94))
    card(s, x, Inches(4.25), Inches(3.72), Inches(1.95), accent=c)
    text(s, Emu(x + Inches(0.36)), Inches(4.55), Inches(3.3), Inches(0.7), [(v, c, True)], size=22, bold=True)
    text(s, Emu(x + Inches(0.36)), Inches(5.4), Inches(3.3), Inches(0.7), l, size=11.5, color=MUT, space=1.12)

# ---------- 招商运营策划及计划 ----------
nxt(); s = slide(); header(s, "05", "OPERATION PLAN", "招商运营策划及计划方案"); footer(s, IDX)
op = [("招商打法", GOLD, ["牌照锚定 + 政策礼包前置，客户「送进来」",
                      "小红书投流引流 → 线下 AI 培训/活动 → 转化租客",
                      "复旦科技园导流 + 高校大厂 + IP/协会资源"]),
      ("运营团队", ACC, ["在地 2 人（招商 + 综合管理，均全方位）",
                     "出海撮合：初期各部门抽调兼职，成熟设全职约 3 人",
                     "挂牌杨浦科技企业服务中心，接政务通道"]),
      ("运营节奏", ACC2, ["月度招商目标 + 去化看板（参考目标，招商不对赌）",
                      "常态化活动完成载体 KPI，运营满年申请补贴（约一二十万）",
                      "企业服务收费体系逐步导入，沉淀会员/私域"])]
for i, (t, c, items) in enumerate(op):
    x = Emu(ML + i * Inches(3.94))
    card(s, x, Inches(1.82), Inches(3.72), Inches(3.05), title=t, items=items, accent=c, bsize=11.3)
rect(s, ML, Inches(5.05), CW, Inches(1.55), grad=[(0, "2A1E55", None), (100, "1E1542", None)], gang=120, line=LINE, radius=True)
rect(s, ML, Inches(5.05), Inches(0.07), Inches(1.55), fill=GOLD)
text(s, Inches(1.1), Inches(5.18), Inches(11), Inches(0.4), "阶段计划（0–3 / 3–6 / 6–12 月）", size=13.5, color=GOLD, bold=True)
bullets(s, Inches(1.1), Inches(5.6), Inches(5.6), Inches(0.9),
        ["0–3 月：定位包装 + 政策梳理 + 首批挂牌 + 样板企业",
         "3–6 月：集中招引 + 首批优惠 + 推介路演 + 政策诊断"], size=11, gap=3, mark=GOLD)
bullets(s, Inches(6.9), Inches(5.6), Inches(5.5), Inches(0.9),
        ["6–12 月：重点补位 + 专精特新/高企培育 + 首场漫展",
         "全周期：活动 IP 化 + 出海撮合 + 补贴反哺 + 复制"], size=11, gap=3, mark=GOLD)

# ---------- 活动安排计划（年度活动日历） ----------
nxt(); s = slide(); header(s, "05", "EVENT CALENDAR", "活动安排计划 · 年度活动日历（约 30 场）"); footer(s, IDX)
table(s, ML, Inches(1.75), CW, [
    ["类型", "频次/年", "内容", "作用"],
    ["AI 培训 / 黑客松", "约 10 场", "OPC 社群、AI 工具链、AIGC 实操", "引流转化 → 租客"],
    ["政策沙龙 / 路演", "约 6 场", "政策解读、专精特新/高企培育、投融资", "政策兑现 + 招商"],
    ["行业对接 / 撮合", "约 6 场", "供应链/IP 出海对接、产业撮合、游学", "出海撮合 + 会员粘性"],
    ["IP / 潮玩活动", "约 4 场", "IP 首发、主题市集、沉浸式展览", "5F 引流 + 商业变现"],
    ["漫展 / 大型活动", "约 2 场", "千人级漫展（门票+赞助）、峰会借势", "声量 + 经营性收入"],
    ["领事专题 / 出海活动", "约 2 场", "国家会客厅领事专题、出海推介", "高端背书 + 出海"]],
    [0.2, 0.13, 0.42, 0.25], sizes=[12, 11, 11, 11], rh=Inches(0.52), head_rh=Inches(0.46), first_col_color=GOLD)
text(s, ML, Inches(6.35), Inches(11.6), Inches(0.3),
     "＊对标魔术空间约 110 场/年：活动场次越多，招商与载体 KPI 越有利；30 场为满足各项要求的基线。", size=9.5, color=SOFT)

# ---------- 18 章节06 排期 ----------
nxt(); divider(IDX, "06", "节点与排期", "MILESTONES & SCHEDULE", "三年合作 · 招商轻运营（不对赌）· 出海为增长引擎")

# ---------- 19 三年排期 ----------
nxt(); s = slide(); header(s, "06", "MILESTONES & SCHEDULE", "三年运营节奏 · 去化与出海双线推进"); footer(s, IDX)
phases = [("第 1 年", "参考 80%", "启动 + 导入", "命名/手册/收费/政策汇编/首批挂牌/样板企业/首批招引 AI·IP·潮玩", ACC),
          ("第 2 年", "参考 85%", "成型 + 提质", "重点客户补位/专精特新·高企培育/企业服务收费/出海撮合起量", ACC2),
          ("第 3 年", "参考 95%", "提升 + 复制", "区级示范点/载体资质/补贴反哺/出海撮合放量/协同周边≤2项目", GOLD)]
rect(s, Inches(1.2), Inches(2.15), Pt(3), Inches(3.9), grad=[(0, H_ACC, None), (100, H_GOLD, None)], gang=90)
y = Inches(2.05)
for nm, tgt, stage, act, c in phases:
    oval(s, Inches(1.05), Emu(y + Inches(0.35)), Inches(0.36), Inches(0.36), fill=c)
    card(s, Inches(1.7), y, Inches(10.78), Inches(1.15), accent=c)
    text(s, Inches(2.0), y, Inches(1.7), Inches(1.15), [(nm, c, True)], size=17, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(3.55), y, Inches(1.7), Inches(1.15), [(tgt, INK, True)], size=16, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(5.35), y, Inches(1.9), Inches(1.15), stage, size=13, color=MUT, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(7.3), y, Inches(5.0), Inches(1.15), act, size=11.5, color=MUT, anchor=MSO_ANCHOR.MIDDLE, space=1.12)
    y = Emu(y + Inches(1.3))

# ---------- 20 章节07 条款报价 ----------
nxt(); divider(IDX, "07", "商业条款与合作报价", "TERMS & QUOTATION", "招商轻运营（不对赌）· 出海重收费 · 联合公司 · 测算")

# ---------- 21 合作模式与原则 ----------
nxt(); s = slide(); header(s, "07", "COOPERATION MODEL", "合作模式与原则 · 轻资产运营"); footer(s, IDX)
mod = [("合作周期", "3 年", "首期三年，满租后享优先续约", GOLD),
       ("协同复制", "≤ 2 个", "可协同周边不超过两个项目，模式可复制", ACC),
       ("合作方式", "轻资产", "招商不对赌；月费转化为出海专项服务费", ACC2),
       ("协议分层", "双协议", "招商轻运营协议 + 出海联合运营协议", GREEN)]
for i, (t, v, d, c) in enumerate(mod):
    x = Emu(ML + i * Inches(2.96))
    card(s, x, Inches(1.85), Inches(2.78), Inches(2.0), accent=c)
    text(s, Emu(x + Inches(0.3)), Inches(2.05), Inches(2.3), Inches(0.4), t, size=13.5, color=MUT, bold=True)
    text(s, Emu(x + Inches(0.3)), Inches(2.5), Inches(2.3), Inches(0.6), [(v, c, True)], size=25, bold=True)
    text(s, Emu(x + Inches(0.3)), Inches(3.2), Inches(2.35), Inches(0.6), d, size=10.5, color=MUT, space=1.14)
text(s, ML, Inches(4.15), Inches(8), Inches(0.4), "合作原则", size=16, color=GOLD, bold=True)
pr = ["空间价值最大化——做产业运营而非二房东租赁",
      "轻资产不对赌——招商板块不背对赌，降低风险",
      "招商佣金强化——按市场 2–3 个月租金取酬",
      "月费转化——原月费并入出海专项服务费（换名）",
      "出海重收费——强化国际出海职能并在此收费",
      "可复制外延——跑通后协同周边项目、拓展新合作"]
for i, p in enumerate(pr):
    x = Emu(ML + (i % 2) * Inches(5.92)); y = Emu(Inches(4.65) + (i // 2) * Inches(0.66))
    rect(s, x, y, Inches(0.1), Inches(0.5), fill=GOLD if i % 2 == 0 else ACC, radius=True)
    text(s, Emu(x + Inches(0.24)), y, Inches(5.6), Inches(0.5), p, size=12, color=INK, anchor=MSO_ANCHOR.MIDDLE)

# ---------- 22 合作报价总表 ----------
nxt(); s = slide(); header(s, "07", "QUOTATION", "合作报价总表 · 招商轻运营（不对赌）+ 出海重收费"); footer(s, IDX)
table(s, ML, Inches(1.72), CW, [
    ["收费项", "标准 / 口径", "支付方", "备注"],
    ["招商佣金（强化）", "成交后 2–3 个月租金（首年不重复）", "合作方", "市场标准，作招商板块主收费；招商不对赌"],
    ["活动执行费", "约 30 场/年 打包 30 万（每场≥30 人）", "合作方", "含策划执行、导师及物料"],
    ["挂牌费", "3 万/个，按数量与影响力分档", "合作方", "国际会议厅费用另计"],
    ["媒体流量费", "10 万/年（含网络投流）", "合作方", "官媒背书 5 万 + 自媒体/投流 5 万"],
    ["出海专项服务费", "出海协议框架内计收（由原月费转化，非月费）", "合作方/合资公司", "含 2–3 人出海团队，强化国际出海职能"],
    ["出海撮合佣金", "平台利润 1–30 万/项目 + 落地佣金 1%–3%", "合资公司", "国际出海收费重点"],
    ["企业增值服务", "工商/政策/知产/财税法务/融资/出海撮合", "入驻企业", "成立联合公司经营分成"],
    ["补贴与资质", "协助申请政府补贴与企业资质", "—", "申请下的补贴归运营平台公司"]],
    [0.17, 0.45, 0.13, 0.25], sizes=[12, 10.5, 10, 10.5], rh=Inches(0.52), head_rh=Inches(0.46), first_col_color=GOLD)

# ---------- 23 媒体报价 + 增值 ----------
nxt(); s = slide(); header(s, "07", "MEDIA & VALUE-ADD", "媒体报价 + 联合公司增值服务"); footer(s, IDX)
rect(s, ML, Inches(1.8), Inches(5.72), Inches(4.75), grad=[(0, "2A1E55", None), (100, "1E1542", None)], gang=120, line=LINE, radius=True)
rect(s, ML, Inches(1.8), Inches(0.07), Inches(4.75), fill=ACC)
text(s, Inches(1.1), Inches(1.98), Inches(5.4), Inches(0.4), [("媒体报价　", INK, True), ("10 万/年（含网络投流）", ACC, True)], size=15.5)
table(s, Inches(1.05), Inches(2.5), Inches(5.4), [
    ["构成", "口径"],
    ["官媒背书 5 万", "党媒/央媒/上海主流媒体宣发，项目背书"],
    ["自媒体+投流 5 万", "小红书投流 + 账号代运营 + 短视频/直播"],
    ["注册自有账号", "园区/IP 官方号，沉淀私域流量"],
    ["活动协同", "活动招募与内容基于活动产出"]],
    [0.32, 0.68], sizes=[11, 10.5], rh=Inches(0.52), head_rh=Inches(0.42), first_col_color=INK)
rect(s, Inches(6.76), Inches(1.8), Inches(5.72), Inches(4.75), grad=[(0, "3A2878", None), (100, "2A1E55", None)], gang=120, line=GOLD, lw=1.25, radius=True)
rect(s, Inches(6.76), Inches(1.8), Inches(0.07), Inches(4.75), fill=GOLD)
text(s, Inches(7.02), Inches(1.98), Inches(5.4), Inches(0.4), "企业增值服务 · 联合公司经营分成", size=15.5, color=GOLD, bold=True)
bullets(s, Inches(7.02), Inches(2.55), Inches(5.3), Inches(2.6), [
    "工商注册、政策申报、知识产权",
    "税收法务、融资顾问、出海撮合",
    "以上通过成立联合公司经营、按约定分成",
    "补贴归属：申请下的补贴归运营平台公司"], size=12.5, gap=11, mark=GOLD)
rect(s, Inches(7.02), Inches(5.25), Inches(5.2), Inches(1.05), grad=[(0, "322462", None), (100, "221848", None)], gang=0, line=LINE, radius=True)
text(s, Inches(7.22), Inches(5.35), Inches(4.9), Inches(0.9),
     [("多元化收入　", GOLD, True), ("联合培训 + 小红书招募增经营性收入；顶层可办漫展/展览（门票+赞助）。", INK, False)], size=11.5, anchor=MSO_ANCHOR.MIDDLE, space=1.15)

# ---------- 媒体合作报价（媒体服务协议） ----------
nxt(); s = slide(); header(s, "07", "MEDIA QUOTATION", "媒体合作报价 · 完整生态传播包"); footer(s, IDX)
table(s, ML, Inches(1.78), Inches(8.35), [
    ["服务项目", "数量", "说明"],
    ["原创内容生产", "3 篇", "结合品牌定位，品牌宣传稿 + 营销活动稿原创采写"],
    ["中央级/全国性/上海主流媒体宣发", "10 篇", "央视/人民网/新华网/第一财经/澎湃/界面等"],
    ["全国社交平台流量投放", "—", "抖音 / 今日头条信息流，投放上海（不含短视频拍摄制作）"],
    ["营销策划与舆情管理咨询", "≤2 次", "营销策划咨询 + 媒体舆情协调处理"]],
    [0.42, 0.13, 0.45], sizes=[12, 11.5, 11], rh=Inches(0.62), head_rh=Inches(0.48), first_col_color=INK)
rect(s, ML, Inches(5.2), Inches(8.35), Inches(1.35), grad=[(0, "3A2878", None), (100, "2A1E55", None)], gang=0, line=GOLD, lw=1.25, radius=True)
text(s, Inches(1.1), Inches(5.35), Inches(8), Inches(0.5),
     [("折后总价　", GOLD, True), ("5 万元", GOLD, True), ("（官媒背书包，曝光 ≥ 150 万；为 10 万/年媒体预算之一）", INK, True)], size=14)
text(s, Inches(1.1), Inches(5.95), Inches(8), Inches(0.5),
     [("交付：", MUT, False), ("合作结束提供完整投放报告（发布链接证明 + 曝光量证明 + 投放分析）", INK, False)], size=11.5)
card(s, Inches(9.42), Inches(1.78), Inches(3.06), Inches(2.35), "服务周期", accent=GOLD,
     body="自签约起至 2026.12.31；服务期满可续，乙方享优先续约权。合作内容可按宣传需要调整。")
card(s, Inches(9.42), Inches(4.28), Inches(3.06), Inches(2.27), "传播策略", accent=ACC, items=[
    "主流媒体引领，权威发声定调",
    "抖音 + 头条：广度覆盖 + 深度触达",
    "信息流精准算法，锁定目标人群"])

# ---------- 24 收益测算 ----------
nxt(); s = slide(); header(s, "07", "ROI", "收益模型 · 对方视角测算"); footer(s, IDX)
kk = [("招商佣金", "2–3 月租金 · 成交计（强化）", GOLD), ("出海专项服务费", "原月费转化 · 出海协议内", ACC),
      ("+ 30 万", "活动打包（约 30 场）", ACC2), ("+ 10 万", "媒体流量费（含投流）", GREEN)]
for i, (v, l, c) in enumerate(kk):
    x = Emu(ML + i * Inches(2.96))
    card(s, x, Inches(1.85), Inches(2.78), Inches(1.5), accent=c)
    text(s, Emu(x + Inches(0.3)), Inches(2.02), Inches(2.3), Inches(0.6), [(v, c, True)], size=21, bold=True)
    text(s, Emu(x + Inches(0.3)), Inches(2.7), Inches(2.35), Inches(0.5), l, size=11, color=MUT, space=1.12)
rect(s, ML, Inches(3.6), Inches(5.72), Inches(2.95), grad=[(0, "2A1E55", None), (100, "1E1542", None)], gang=120, line=LINE, radius=True)
rect(s, ML, Inches(3.6), Inches(0.07), Inches(2.95), fill=GOLD)
text(s, Inches(1.1), Inches(3.78), Inches(5.4), Inches(0.4), "对方视角", size=14.5, color=INK, bold=True)
table(s, Inches(1.05), Inches(4.25), Inches(5.35), [
    ["维度", "口径"], ["对方空间成本", "≈554 万/年（2.3 元/㎡/天）"],
    ["招商板块", "佣金 2–3 月 + 活动 30 万 + 媒体 10 万（不对赌）"], ["出海板块", "出海专项服务费 + 撮合佣金（重收费）"]],
    [0.28, 0.72], sizes=[11.5, 10.5], rh=Inches(0.5), head_rh=Inches(0.44), first_col_color=ACC)
card(s, Inches(6.76), Inches(3.6), Inches(5.72), Inches(2.95), "为什么算得过账", accent=GREEN, items=[
    "招商轻资产不对赌，降低双方风险",
    "佣金强化至 2–3 月，贴合市场标准",
    "月费转出海专项服务费，出海为增长引擎",
    "出海撮合 + 联合公司 + 国家会客厅第二曲线"])

# ---------- 章节08 出海与国际会客厅 ----------
nxt(); divider(IDX, "08", "出海撮合与国际会客厅", "GLOBAL & CONSUL SALON", "国家会客厅 · 领事资源 · 出海撮合业务闭环")

# ---------- 国际会议厅 · 国家会客厅（匿名） ----------
nxt(); s = slide(); header(s, "08", "NATIONAL SALON", "国际会议厅 · 「国家会客厅」（旗舰载体）"); footer(s, IDX)
rect(s, ML, Inches(1.72), CW, Inches(0.92), grad=[(0, "3A2878", None), (100, "2A1E55", None)], gang=0, line=GOLD, lw=1.25, radius=True)
text(s, Inches(1.1), Inches(1.72), Inches(11.2), Inches(0.92),
     [("载体定位　", GOLD, True), ("联合国际领事资源合作方（驻沪数十国总领事网络与领团背书），依托上海核心区地标高层空间，设「国家会客厅」——外事接待 · 文化展示 · 招商对接 · 企业服务四位一体。", INK, False)],
     size=13, anchor=MSO_ANCHOR.MIDDLE, space=1.18)
tri = [("三方协同", GOLD, ["复旦政策研究中心：智库背书 / 课题 / 培训",
                       "杨浦科企联：企业会员 / 获客 / 招商落地",
                       "领事资源合作方：领事网络 / 空间 / 活动落地"]),
       ("六大产品线", ACC, ["国家会客厅（旗舰，冠名 + 会籍）",
                        "领事专题活动（小 / 中 / 大三档）",
                        "高端商务出海深度游 · 企业出海对接",
                        "会员会籍 · 文化传播 / 智库 / 培训"]),
       ("收入来源", ACC2, ["国家冠名 80–300 万 / 联合冠名 30–80 万",
                       "企业常驻会籍 10–30 万 / 单场地 + 领事出席 3–10 万",
                       "园区·政府招商包 30–100 万 / 项目",
                       "出海：单国服务包 20–50 万 / 落地佣金 1–5%"])]
for i, (t, c, items) in enumerate(tri):
    x = Emu(ML + i * Inches(3.94))
    card(s, x, Inches(2.85), Inches(3.72), Inches(3.35), title=t, items=items, accent=c, bsize=11)
text(s, ML, Inches(6.35), Inches(11.6), Inches(0.3),
     "＊为商务保密，本页对合作方名称作模糊化处理；具体主体、空间与权益以正式协议为准。", size=9.5, color=SOFT)

# ---------- 出海撮合业务合作模式 ----------
nxt(); s = slide(); header(s, "08", "GO-GLOBAL MODEL", "出海撮合业务 · 合作模式与闭环"); footer(s, IDX)
flow2 = [("展厅免租", "5F 展厅对优质企业免租，约定租金标准 + 物业费", GOLD),
         ("订单扣点", "出海达成订单后按扣点反哺租金", ACC),
         ("超额分成", "超出租金部分进入合资公司收益", ACC2),
         ("合资分配", "运营平台方与合作方按约定分配（平台方约 30%）", GREEN)]
for i, (t, d, c) in enumerate(flow2):
    x = Emu(ML + i * Inches(2.98))
    card(s, x, Inches(1.78), Inches(2.8), Inches(1.75), title=t, body=d, accent=c, tsize=14.5, bsize=10.8)
    if i < 3:
        chev(s, Emu(x + Inches(2.8)), Emu(Inches(1.78) + Inches(0.5)))
card(s, ML, Inches(3.75), Inches(5.72), Inches(2.6), "收入模式", accent=GOLD, items=[
    "出海撮合平台利润：单项目 1–30 万",
    "落地成功佣金：项目额 1%–3%（按行业浮动）",
    "展厅订单扣点反哺租金 + 超额进合资公司",
    "产业撮合会员体系：年度顾问 / 游学 / 高级专家",
    "国家会客厅：冠名 / 会籍 / 园区政府招商包"])
card(s, Inches(6.76), Inches(3.75), Inches(5.72), Inches(2.6), "落地要点", accent=ACC, items=[
    "行业聚焦：供应链（玩具/毛绒/潮玩 OEM→品牌）+ 内容（AI 短剧/IP 授权）",
    "对接层级：连领事 → 对接所在国商务部/部长 → 洽谈落地",
    "团队：初期各部门抽调兼职（1 人起），成熟设全职（约 3 人），单独成立公司",
    "沪穗联动：与广州自贸港/电商出海项目协同，两地共享资源"])

# ---------- 出海全链条 · 流程与服务（运营框架） ----------
nxt(); s = slide(); header(s, "08", "GO-GLOBAL FULL CHAIN", "出海全链条 · 流程与服务（运营框架）"); footer(s, IDX)
steps6 = [("01", "集聚入驻", "3F 孵化办公\n5F 展贸集群", GOLD),
          ("02", "内容与产品", "AIGC/AI 短剧制作\nIP 授权 · 选品打样", ACC),
          ("03", "海外通路", "TikTok/领英/短剧\n跨境电商 · 展销", ACC2),
          ("04", "订单撮合", "5F 展厅接单\n领事/商务对接", GREEN),
          ("05", "履约交付", "报关物流 · 结汇\n财税合规", ROSE),
          ("06", "增长复购", "品牌沉淀 · 复投\n反哺租金/合资收益", GOLD)]
cwg = Inches(1.86)
for i, (n, t, d, c) in enumerate(steps6):
    x = Emu(ML + i * (cwg + Inches(0.1)))
    card(s, x, Inches(1.75), cwg, Inches(2.0), accent=c)
    text(s, Emu(x + Inches(0.22)), Inches(1.9), Inches(1.2), Inches(0.5), [(n, c, True, FONT_EN)], size=20, font=FONT_EN)
    text(s, Emu(x + Inches(0.22)), Inches(2.38), Emu(cwg - Inches(0.35)), Inches(0.4), t, size=12.5, color=INK, bold=True)
    text(s, Emu(x + Inches(0.22)), Inches(2.78), Emu(cwg - Inches(0.35)), Inches(0.9), d, size=9.5, color=MUT, space=1.15)
srv = [("内容出海服务链", ACC, ["AI 短剧：剧本→AIGC 制作→翻译配音→TikTok/海外短剧平台分发→广告/付费分成",
                        "数字 IP：国漫 IP 海外独家授权 → 海外手办/玩具开发 → 跨境渠道销售",
                        "领英出海：LinkedIn B2B 海外获客 / 品牌出海 / 决策人触达"]),
       ("货品出海服务链", GOLD, ["供应链：玩具/毛绒/潮玩 OEM → 设计赋能 → 品牌化 → 海外订单",
                        "通路：TikTok / 领英 / 跨境电商（亚马逊·TikTok Shop）+ 海外展销 + 自贸港",
                        "配套：报关物流 · 出口退税 · 结汇 · 财税法务（联合公司承接）"])]
for i, (t, c, items) in enumerate(srv):
    x = Emu(ML + i * Inches(5.92))
    card(s, x, Inches(4.0), Inches(5.72), Inches(2.35), title=t, items=items, accent=c, bsize=11)
text(s, ML, Inches(6.5), Inches(11.6), Inches(0.3),
     "＊全链条服务由运营平台 + 联合公司 + 出海合作资源共同交付；收益按撮合佣金/分成计入合资公司。", size=9.5, color=SOFT)

# ---------- 领事馆路径 · 海外订单达成逻辑 ----------
nxt(); s = slide(); header(s, "08", "CONSULATE PATH", "领事馆路径 · 产业集聚 → 海外订单达成逻辑"); footer(s, IDX)
steps5 = [("STEP 1", "产业集聚", "各类产业企业集聚平台：\n供应链 / IP 内容 / AI 应用\n形成可展示的产品池", GOLD),
          ("STEP 2", "需求画像", "按国别梳理供需清单：\n企业出海意向 × 国别\n市场准入与政策研判", ACC),
          ("STEP 3", "领事撮合", "国家会客厅 / 领事专题活动：\n连领事 → 对接所在国\n商务部 / 部长级洽谈", ACC2),
          ("STEP 4", "商务对接", "国别推介会 / 买家团来访 /\n出海考察团（B2B 配对）\n展厅选品 · 样品确认", GREEN),
          ("STEP 5", "订单落地", "签约 → 履约交付 → 结算\n订单扣点反哺租金\n超额进合资公司分配", ROSE)]
cwp = Inches(2.24)
for i, (n, t, d, c) in enumerate(steps5):
    x = Emu(ML + i * (cwp + Inches(0.11)))
    card(s, x, Inches(1.75), cwp, Inches(2.6), accent=c)
    text(s, Emu(x + Inches(0.24)), Inches(1.9), Inches(1.6), Inches(0.4), [(n, c, True, FONT_EN)], size=13, font=FONT_EN)
    text(s, Emu(x + Inches(0.24)), Inches(2.28), Emu(cwp - Inches(0.4)), Inches(0.4), t, size=14, color=INK, bold=True)
    text(s, Emu(x + Inches(0.24)), Inches(2.72), Emu(cwp - Inches(0.4)), Inches(1.5), d, size=9.8, color=MUT, space=1.18)
    if i < 4:
        chev(s, Emu(x + cwp), Inches(2.7))
rect(s, ML, Inches(4.6), CW, Inches(1.85), grad=[(0, "3A2878", None), (100, "2A1E55", None)], gang=0, line=GOLD, lw=1.25, radius=True)
text(s, Inches(1.1), Inches(4.72), Inches(11.2), Inches(0.4),
     [("双向撮合方法　", GOLD, True), ("「走出去」与「引进来」同一套领事通路，双向取酬。", INK, False)], size=13)
bullets(s, Inches(1.1), Inches(5.15), Inches(5.6), Inches(1.2),
        ["走出去：企业产品池 → 国别推介 → 海外买家 → 订单出口",
         "引进来：领事引荐海外品牌/买家 → 5F 展厅落地 → 采购中国供应链"], size=11, gap=5, mark=GOLD)
bullets(s, Inches(6.9), Inches(5.15), Inches(5.5), Inches(1.2),
        ["高频活动养通路：领事午餐会/国别论坛/买家团（纳入年度 30 场）",
         "取酬：撮合佣金 1%–3% + 单国服务包 + 会籍 + 订单扣点反哺租金"], size=11, gap=5, mark=GOLD)

# ---------- 大湾区 · 广州出海资源叠加 ----------
nxt(); s = slide(); header(s, "08", "GBA RESOURCES", "大湾区出海资源叠加 · 复旦系 + 高校（沪穗联动）"); footer(s, IDX)
gba = [("复旦珠海创新研究院", GOLD, "九大公共平台（集成电路/大数据智算/工业仿真/海洋科技/物联网智慧城市/先进材料）；「手心研发·手背转化」；服务横琴粤澳深合区"),
       ("葡语国家科技合作平台", ACC, "横琴对接葡语系（巴西/葡萄牙/安哥拉/莫桑比克）唯一国家级平台；市场调研·本地匹配·海外展会·合规咨询"),
       ("珠澳双循环出海窗口", ACC2, "澳门自由港 + 葡语商贸平台；「内地研发 + 澳门展销 + 葡语国家落地」；设海外总部/海外仓中转，规避关税壁垒"),
       ("粤港澳复旦成果转化中心", GREEN, "复旦技术转移落地南沙；技术授权海外/国际专利/海外临床与认证；港澳跨境通道 + 粤科复旦母基金等出海基金矩阵"),
       ("大湾区精准医学研究院", ROSE, "省市南沙复旦五方共建；南沙唯一万笼级实验动物中心；肿瘤/代谢/生殖精准诊疗、类器官、医疗器械转化"),
       ("深圳复旦研究院", ACC, "河套深港科创合作区；AI/前沿医学/集成电路；联动香港，打通深港产学研")]
for i, (t, c, d) in enumerate(gba):
    x = Emu(ML + (i % 3) * Inches(3.94)); y = Emu(Inches(1.78) + (i // 3) * Inches(1.72))
    card(s, x, y, Inches(3.72), Inches(1.55), title=t, body=d, accent=c, tsize=13, bsize=10)
rect(s, ML, Inches(5.28), CW, Inches(1.32), grad=[(0, "3A2878", None), (100, "2A1E55", None)], gang=0, line=GOLD, lw=1, radius=True)
text(s, Inches(1.1), Inches(5.4), Inches(11.2), Inches(0.4), [("高校出海资源　", GOLD, True), ("暨南大学 · 华侨大学 · 广东外语外贸大学 · 广州大学黄埔研究院 + 复旦广州校友会", INK, True)], size=12.5)
bullets(s, Inches(1.1), Inches(5.85), Inches(11), Inches(0.7),
        ["广外＝出海第一梯队：小语种（东盟/葡语/西语）· 海丝国别智库 · 跨境合规 · 国际贸易 · 海外市场调研",
         "广州大学黄埔研究院：数字经济 · 跨境新媒体 · 海外独立站 · 智能制造国际转化"], size=10.5, gap=3, mark=GOLD)

# ---------- 战略资源探讨 · 东方枢纽 + 谷歌 + 领英 ----------
nxt(); s = slide(); header(s, "08", "STRATEGIC EXPLORATION", "战略资源探讨 · 东方枢纽 / 谷歌 / 领英"); footer(s, IDX)
ex = [("东方枢纽", GOLD, ["国际商务合作区 + 综保区 + 空铁枢纽（临港/浦东方向）",
                     "探讨：出海企业「境内关外」通道、保税展示、国际中转",
                     "作用：为出海订单与海外仓提供口岸与政策承接点"]),
      ("谷歌 Google", ACC, ["探讨：Google Ads / YouTube 出海投放与海外获客",
                        "探讨：谷歌在领事馆/国别活动中的技术与流量角色",
                        "作用：叠加海外数字营销通路，服务内容与货品出海"]),
      ("领英 LinkedIn", ACC2, ["B2B 海外决策人触达、品牌出海、海外招募",
                          "与展会/领事撮合互补，沉淀海外客户资产",
                          "作用：国际出海职能的线上主通路之一"])]
for i, (t, c, items) in enumerate(ex):
    x = Emu(ML + i * Inches(3.94))
    card(s, x, Inches(1.8), Inches(3.72), Inches(3.05), title=t, items=items, accent=c, bsize=11)
rect(s, ML, Inches(5.05), CW, Inches(1.5), grad=[(0, "2A1E55", None), (100, "1E1542", None)], gang=120, line=GOLD, lw=1, radius=True)
rect(s, ML, Inches(5.05), Inches(0.07), Inches(1.5), fill=GOLD)
text(s, Inches(1.1), Inches(5.18), Inches(11), Inches(0.4), [("说明　", GOLD, True), ("本页为战略资源探讨方向（待落实），用以强化国际出海职能并在出海板块形成收费。", INK, False)], size=12)
bullets(s, Inches(1.1), Inches(5.62), Inches(11), Inches(0.8),
        ["收费落点：出海专项服务费 + 撮合佣金 + 海外营销投放代运营 + 国家会客厅冠名/会籍",
         "东方枢纽与谷歌合作以实际对接结果为准，不构成承诺"], size=10.5, gap=3, mark=GOLD)

# ========== 三大合作展开（谷歌 / 领英 / 东方枢纽）==========
def collab(idx, eyebrow, title, lead, content, fees, costs, land, status, scolor):
    s = slide(); header(s, "08", eyebrow, title); footer(s, idx)
    rect(s, ML, Inches(1.68), CW, Inches(0.92), grad=[(0, "3A2878", None), (100, "2A1E55", None)], gang=0, line=GOLD, lw=1.25, radius=True)
    text(s, Inches(1.1), Inches(1.68), Inches(11.2), Inches(0.92),
         [("合作模式　", GOLD, True), (lead, INK, False)], size=12.5, anchor=MSO_ANCHOR.MIDDLE, space=1.18)
    card(s, ML, Inches(2.78), Inches(3.72), Inches(2.5), "合作内容", items=content, accent=ACC, bsize=11)
    card(s, Inches(4.79), Inches(2.78), Inches(3.72), Inches(2.5), "收费（含单位）", items=fees, accent=GOLD, bsize=11)
    card(s, Inches(8.78), Inches(2.78), Inches(3.7), Inches(2.5), "成本（含单位）", items=costs, accent=ROSE, bsize=11)
    rect(s, ML, Inches(5.45), CW, Inches(1.1), grad=[(0, "2A1E55", None), (100, "1E1542", None)], gang=120, line=LINE, radius=True)
    rect(s, ML, Inches(5.45), Inches(0.07), Inches(1.1), fill=scolor)
    text(s, Inches(1.1), Inches(5.56), Inches(11.2), Inches(0.4), [("收费落点　", GOLD, True), (land, INK, False)], size=11.5)
    text(s, Inches(1.1), Inches(6.02), Inches(11.2), Inches(0.4), [("状态　", scolor, True), (status, MUT, False)], size=11)
    return s

nxt(); collab(IDX, "GOOGLE PARTNERSHIP", "谷歌合作 · 合作模式展开",
    "以园区/平台名义对接谷歌出海生态（Google Ads / YouTube / Google Cloud），为入驻企业提供海外投放开户、代运营与出海营销培训；并探讨谷歌在领事馆/国别活动中的技术与流量支持角色。",
    ["海外广告开户、充值与合规", "YouTube / 搜索出海投放代运营", "出海营销培训与认证", "谷歌出海资源与国别活动对接"],
    ["投放代运营 1–3 万/月/家", "或按投放额 10%–15% 计", "开户返点分成（平台侧）", "培训 按场/人计"],
    ["代运营人力（兼职分摊）", "投放媒介费由客户承担", "账户合规/风控成本", "—"],
    "出海专项服务费 + 海外营销代运营 + 返点分成。",
    "探讨中 / 待落实，以实际对接结果为准，不构成承诺。", ACC)

nxt(); collab(IDX, "LINKEDIN PARTNERSHIP", "领英合作 · 合作模式展开",
    "以平台名义对接领英营销解决方案，为出海企业做 B2B 品牌页、广告投放、决策人触达与海外招募；纳入出海全链条线上主通路，与展会/领事撮合互补。",
    ["企业领英账号搭建与代运营", "B2B 广告投放 / 决策人触达", "海外 KOL / 人才招募", "领英出海培训"],
    ["账号搭建 0.5–1 万/家（一次性）", "代运营 1–3 万/月/家", "或按投放额 10%–15% 计", "培训 按场计"],
    ["代运营人力（兼职分摊）", "投放媒介费由客户承担", "领英资源合作分成", "—"],
    "出海专项服务费 + 代运营分成。",
    "探讨中 / 待落实，沉淀海外客户资产为长期价值。", ACC2)

nxt(); collab(IDX, "EAST-HUB PARTNERSHIP", "东方枢纽合作 · 合作模式展开",
    "对接东方枢纽（浦东综保区 / 国际商务合作区 / 空铁枢纽），为出海企业提供境内关外通道、保税展示、海外仓与国际中转、口岸与政策承接，作为出海物流与合规落地节点。",
    ["保税展示 / 展销", "海外仓 / 国际中转", "跨境电商 / 出口退税通道", "口岸政策承接与合规"],
    ["服务对接费（待定）", "撮合佣金（待定）", "仓储 / 中转分成（待定）", "—"],
    ["合作模式未明，成本待定", "—", "—", "—"],
    "若落地：服务对接费 + 撮合佣金 + 仓储分成。",
    "中保区（东方枢纽）合作模式尚未明确，暂不纳入本次预算；作为加分项后补（主干先行·枝叶后补）。", ROSE)

# ---------- 出海费用清单（轻运营 · 挂牌/代运营/资源导入） ----------
nxt(); s = slide(); header(s, "08", "GO-GLOBAL FEES", "出海费用清单 · 挂牌 / 代运营 / 资源导入（轻运营）"); footer(s, IDX)
rect(s, ML, Inches(1.68), CW, Inches(0.82), grad=[(0, "3A2878", None), (100, "2A1E55", None)], gang=0, line=GOLD, lw=1.25, radius=True)
text(s, Inches(1.1), Inches(1.68), Inches(11.2), Inches(0.82),
     [("方式　", GOLD, True), ("出海板块以「轻运营」推进：以运营活动 + 资源对接为主，不重投入；活动场次或按招人方式后续再定。以下为出海成本项（供甲方预算）。", INK, False)],
     size=12.5, anchor=MSO_ANCHOR.MIDDLE, space=1.16)
table(s, ML, Inches(2.72), CW, [
    ["费用项", "内容", "金额（含单位）"],
    ["① 出海挂牌费", "国家 / 北欧会客厅挂牌 + 领事资源背书 + 出海身份/牌照", "10–20 万元（一次性 / 按牌）"],
    ["② 代运营费用", "领英 / 谷歌 / 短剧 / 社媒 海外账号搭建与投放代运营（媒介投放费由企业承担）", "打包 15–30 万/年（或 1–3 万/月/家）"],
    ["③ 资源导入费用", "大湾区 / 复旦系 / 高校 / 领事 / 买家团 / IP 资源对接与导入", "打包 10–20 万/年（或按项目 1–5 万/次）"]],
    [0.2, 0.54, 0.26], sizes=[12, 10.5, 11], rh=Inches(0.72), head_rh=Inches(0.48), first_col_color=GOLD)
rect(s, ML, Inches(5.5), CW, Inches(1.05), grad=[(0, "2A1E55", None), (100, "1E1542", None)], gang=120, line=GOLD, lw=1, radius=True)
rect(s, ML, Inches(5.5), Inches(0.07), Inches(1.05), fill=GOLD)
text(s, Inches(1.1), Inches(5.6), Inches(11.2), Inches(0.4),
     [("费用合计（示意）　", GOLD, True), ("一次性：挂牌 10–20 万；年度：代运营 + 资源导入 约 25–50 万/年 → 首年约 35–70 万。", INK, False)], size=12)
text(s, Inches(1.1), Inches(6.04), Inches(11.2), Inches(0.4),
     [("备注　", MUT, True), ("不含出海撮合佣金（平台利润 1–30 万/项目 + 落地 1%–3%，按成交另计）与海外媒介投放费（企业承担）；活动/招人方式待定后细化。", MUT, False)], size=10.5, space=1.12)

# ---------- 营收逻辑 · 收入结构（3F+5F 整体） ----------
nxt(); s = slide(); header(s, "08", "REVENUE MODEL", "营收逻辑 · 收入结构（3F + 5F 整体）"); footer(s, IDX)
rect(s, ML, Inches(1.68), CW, Inches(0.82), grad=[(0, "3A2878", None), (100, "2A1E55", None)], gang=0, line=GOLD, lw=1.25, radius=True)
text(s, Inches(1.1), Inches(1.68), Inches(11.2), Inches(0.82),
     [("收入逻辑　", GOLD, True), ("园区空间盘（租金/物业）＝出海反哺来源；运营平台收入＝招商佣金（强化）+ 出海专项服务费 + 撮合分成 + 经营性/补贴。", INK, False)],
     size=12, anchor=MSO_ANCHOR.MIDDLE, space=1.15)
# 左：园区空间盘（背景）
rect(s, ML, Inches(2.66), Inches(5.72), Inches(3.7), grad=[(0, "2A1E55", None), (100, "1E1542", None)], gang=120, line=LINE, radius=True)
rect(s, ML, Inches(2.66), Inches(0.07), Inches(3.7), fill=ACC)
text(s, Inches(1.08), Inches(2.8), Inches(5.4), Inches(0.4), "① 园区空间盘（甲方口径 · 出海反哺来源）", size=13, color=ACC, bold=True)
table(s, Inches(1.03), Inches(3.3), Inches(5.4), [
    ["科目", "年化（示意）"],
    ["3F 办公租金 2850㎡", "约 400 万（85% 去化 · ~4.5 元/㎡/天）"],
    ["5F 展贸租金 + 物业", "约 300 万（部分免租 · 订单扣点反哺）"],
    ["6600㎡ 物业费", "约 200 万（1.0 元/㎡/天 计）"],
    ["空间盘合计", "约 900 万/年（满租稳定期）"]],
    [0.46, 0.54], sizes=[10.5, 10.5], rh=Inches(0.5), head_rh=Inches(0.42), first_col_color=INK)
# 右：运营平台收入（乙方口径）
rect(s, Inches(6.76), Inches(2.66), Inches(5.72), Inches(3.7), grad=[(0, "2A1E55", None), (100, "1E1542", None)], gang=120, line=GOLD, lw=1.25, radius=True)
rect(s, Inches(6.76), Inches(2.66), Inches(0.07), Inches(3.7), fill=GOLD)
text(s, Inches(7.04), Inches(2.8), Inches(5.4), Inches(0.4), "② 运营平台收入（乙方口径）", size=13, color=GOLD, bold=True)
table(s, Inches(6.99), Inches(3.3), Inches(5.4), [
    ["收入类别", "年化（示意）"],
    ["招商佣金（强化）", "2–3 月租金/成交（约 60–120 万）"],
    ["活动 + 媒体", "活动 30 万 + 媒体 10 万 + 挂牌"],
    ["出海专项服务费", "原月费转化，出海协议内计收"],
    ["出海撮合 + 增值分成", "约 150–450 万（联合公司）"]],
    [0.4, 0.6], sizes=[10.5, 10.5], rh=Inches(0.5), head_rh=Inches(0.42), first_col_color=INK)
text(s, ML, Inches(6.5), Inches(11.6), Inches(0.3),
     "＊示意测算，用于沟通商业空间；租金/去化/分成以物业条件、签约与执行节奏为准。运营平台稳定期年收入区间约 380–740 万。", size=9.5, color=SOFT)

# ---------- 营收逻辑 · 成本与盈亏测算 ----------
nxt(); s = slide(); header(s, "08", "REVENUE MODEL", "营收逻辑 · 成本结构与盈亏测算"); footer(s, IDX)
rect(s, ML, Inches(1.7), Inches(5.72), Inches(3.55), grad=[(0, "2A1E55", None), (100, "1E1542", None)], gang=120, line=LINE, radius=True)
rect(s, ML, Inches(1.7), Inches(0.07), Inches(3.55), fill=ROSE)
text(s, Inches(1.08), Inches(1.84), Inches(5.4), Inches(0.4), "成本结构（运营平台 · 年化）", size=13, color=ROSE, bold=True)
table(s, Inches(1.03), Inches(2.34), Inches(5.4), [
    ["成本科目", "年化（示意）"],
    ["团队人力（在地 2 + 出海 1–3 + 兼职）", "100–180 万"],
    ["活动执行（约 30 场）", "约 30 万"],
    ["媒体 + 网络投流", "约 10 万"],
    ["运营杂费 / 差旅 / 办公", "20–40 万"],
    ["硬件装修摊销", "甲方为主，乙方轻资产"],
    ["成本合计", "约 160–260 万/年"]],
    [0.62, 0.38], sizes=[10, 10.5], rh=Inches(0.44), head_rh=Inches(0.4), first_col_color=INK)
# 三情景
sc = [("保守", "收入 ~380 万", "成本 ~185 万", "净 +195 万", ACC),
      ("中性", "收入 ~520 万", "成本 ~220 万", "净 +300 万", GOLD),
      ("乐观", "收入 ~740 万", "成本 ~260 万", "净 +480 万", GREEN)]
for i, (t, r, c0, n, c) in enumerate(sc):
    y = Emu(Inches(1.7) + i * Inches(0.9))
    rect(s, Inches(6.76), y, Inches(5.72), Inches(0.78), grad=[(0, "2A1E55", None), (100, "1E1542", None)], gang=120, line=LINE, radius=True)
    rect(s, Inches(6.76), y, Inches(0.07), Inches(0.78), fill=c)
    text(s, Inches(7.02), y, Inches(1.1), Inches(0.78), [(t, c, True)], size=15, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(8.0), y, Inches(2.0), Inches(0.78), r, size=11.5, color=MUT, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(9.5), y, Inches(1.6), Inches(0.78), c0, size=11.5, color=MUT, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(11.0), y, Inches(1.5), Inches(0.78), [(n, c, True)], size=13, anchor=MSO_ANCHOR.MIDDLE)
rect(s, ML, Inches(5.45), CW, Inches(1.1), grad=[(0, "322462", None), (100, "221848", None)], gang=0, line=GOLD, lw=1, radius=True)
text(s, Inches(1.1), Inches(5.56), Inches(11.2), Inches(0.4), [("盈亏与现金流　", GOLD, True), ("招商轻资产不对赌、佣金强化保基本盘；出海专项服务费 + 撮合为最大增长弹性，三情景净收益均为正。", INK, False)], size=12)
bullets(s, Inches(1.1), Inches(6.0), Inches(11), Inches(0.5),
        ["轻资产：装修硬件以甲方投入为主，乙方前期投入低、回正快；补贴与增值/出海形成第二增长曲线"], size=11, gap=2, mark=GOLD)
text(s, ML, Inches(6.62), Inches(11.6), Inches(0.28), "＊三情景为示意测算区间，实际以签约、去化与出海落地为准。", size=9, color=SOFT)

# ---------- 财务测算 · 政策申报相关收入（能否拿到现金） ----------
nxt(); s = slide(); header(s, "08", "POLICY INCOME", "财务测算 · 政策申报收入（能否拿到现金）"); footer(s, IDX)
table(s, ML, Inches(1.7), CW, [
    ["政策", "金额（含单位）", "平台可实现现金（单位）", "现金确定性"],
    ["载体运营/活动补贴（科技服务业·成果转化平台）", "认定后≤100 万/年；运营满年约 10–20 万/年", "直接到平台 10–20 万/年", "现金·较确定（需认定+满年）"],
    ["高新技术企业认定", "首次 20 万/次（企业·一次性）", "申报服务费 2–8 万/家", "现金·确定（服务费）"],
    ["专精特新", "市级 10 / 国家小巨人 30 万（企业·一次性）", "服务费 2–8 万/家", "现金·确定（服务费）"],
    ["科技型中小企业备案", "资格（企业·无直接现金）", "服务费 0.3–2 万/家", "现金·确定（服务费）"],
    ["AI/大数据房租补贴", "2 元/天/㎡·≤100 万/年（企业·3 年）", "代申请成功费 到账 5%–15%", "现金·依到账"],
    ["政府补贴项目（各类）", "视项目（企业）", "成功费 到账额 5%–15%", "现金·以兑现为准"],
    ["科技创新券", "企业≤30 万·团队≤10 万/年（抵付服务）", "作服务机构核销分成", "现金·需入驻创新券平台"],
    ["算力补贴 + 云折扣", "最高 50%（企业）", "— 招商抓手（非平台现金）", "非现金（招商用）"]],
    [0.32, 0.28, 0.24, 0.16], sizes=[10, 9.5, 9.5, 9.5], rh=Inches(0.44), head_rh=Inches(0.42), first_col_color=GOLD)
rect(s, ML, Inches(5.95), CW, Inches(0.72), grad=[(0, "3A2878", None), (100, "2A1E55", None)], gang=0, line=GOLD, lw=1, radius=True)
text(s, Inches(1.1), Inches(5.95), Inches(11.2), Inches(0.72),
     [("关键提示　", GOLD, True), ("计入预算的政策收入＝平台可实现现金（申报服务费 + 成功分成 + 载体运营补贴）；产业补贴主体流向企业/载体，平台按约定分成或成功费取酬。未来预判：微短剧专项、模塑申城三券、YOUNG立方直播券/活动补贴（窗口开放后追加）。", INK, False)],
     size=10.5, anchor=MSO_ANCHOR.MIDDLE, space=1.15)

# ---------- 财务测算 · 国际会客厅（北欧会客厅） ----------
nxt(); s = slide(); header(s, "08", "SALON P&L", "财务测算 · 国际会客厅（北欧会客厅）"); footer(s, IDX)
rect(s, ML, Inches(1.68), CW, Inches(0.98), grad=[(0, "3A2878", None), (100, "2A1E55", None)], gang=0, line=GOLD, lw=1.25, radius=True)
text(s, Inches(1.1), Inches(1.68), Inches(11.2), Inches(0.98),
     [("合作模式　", GOLD, True), ("联合国际领事资源合作方，以「北欧/国别会客厅」为载体，在核心区高层空间做外事接待·文化展示·招商对接·企业出海；平台负责运营策划、活动执行与出海撮合，按「冠名 + 会籍 + 活动 + 出海对接」收费，与领事资源合作方分润。", INK, False)],
     size=12, anchor=MSO_ANCHOR.MIDDLE, space=1.16)
rect(s, ML, Inches(2.82), Inches(5.72), Inches(2.5), grad=[(0, "2A1E55", None), (100, "1E1542", None)], gang=120, line=LINE, radius=True)
rect(s, ML, Inches(2.82), Inches(0.07), Inches(2.5), fill=GOLD)
text(s, Inches(1.08), Inches(2.95), Inches(5.4), Inches(0.4), "收费（合理·不高估）", size=13, color=GOLD, bold=True)
table(s, Inches(1.03), Inches(3.45), Inches(5.4), [
    ["项目", "金额（单位）"],
    ["冠名/联合冠名", "30–80 万/年"],
    ["企业会籍", "10 万/家/年（保守 3–5 家）"],
    ["单场活动（领事出席）", "3–8 万/场（6–8 场/年）"],
    ["出海对接服务", "单国服务包 20–50 万/项目"]],
    [0.42, 0.58], sizes=[10, 10], rh=Inches(0.42), head_rh=Inches(0.38), first_col_color=INK)
rect(s, Inches(6.76), Inches(2.82), Inches(5.72), Inches(2.5), grad=[(0, "2A1E55", None), (100, "1E1542", None)], gang=120, line=LINE, radius=True)
rect(s, Inches(6.76), Inches(2.82), Inches(0.07), Inches(2.5), fill=ROSE)
text(s, Inches(7.02), Inches(2.95), Inches(5.4), Inches(0.4), "成本（含单位）", size=13, color=ROSE, bold=True)
table(s, Inches(6.99), Inches(3.45), Inches(5.4), [
    ["项目", "金额（单位）"],
    ["活动执行", "3–5 万/场"],
    ["领事/嘉宾对接接待", "2–5 万/场"],
    ["空间/场地分摊", "按次或年（合作方空间）"],
    ["合作方分润", "领事资源方 40%–50%"]],
    [0.42, 0.58], sizes=[10, 10], rh=Inches(0.42), head_rh=Inches(0.38), first_col_color=INK)
rect(s, ML, Inches(5.48), CW, Inches(1.05), grad=[(0, "322462", None), (100, "221848", None)], gang=0, line=GOLD, lw=1, radius=True)
text(s, Inches(1.1), Inches(5.58), Inches(11.2), Inches(0.9),
     [("合理预估（不高估）　", GOLD, True), ("首年收入 ~80–120 万，平台净 ~20–40 万/年；成熟期收入 ~150–250 万，平台净 ~50–90 万/年（分润与成本后）。数据为保守测算，待与合作方最终确认。", INK, False)],
     size=11.5, anchor=MSO_ANCHOR.MIDDLE, space=1.15)

# ---------- 财务测算 · 企业服务（项目/收费/成本） ----------
nxt(); s = slide(); header(s, "08", "ENTERPRISE SERVICE P&L", "财务测算 · 企业服务（项目 / 收费 / 成本）"); footer(s, IDX)
table(s, ML, Inches(1.7), CW, [
    ["服务项目", "收费（含单位）", "成本（含单位）"],
    ["工商注册 · 代账", "注册引流免费；代账 3000–6000 元/年/家", "第三方分成 + 人力"],
    ["政策申报", "普通 0.3–2 万/项；高企/专精特新 2–8 万/项；补贴成功费 到账 5%–15%", "人力 + 材料杂费"],
    ["知识产权", "商标 800–2000 元/件；软著 2000–5000 元/件；专利代理 5000–1 万/件", "第三方分成"],
    ["财税 · 法务", "顾问 1–5 万/年/家", "第三方分成"],
    ["融资顾问", "成功佣金 融资额 1%–3%", "人力"],
    ["出海撮合", "平台利润 1–30 万/项目 + 落地佣金 1%–3%", "出海团队（单列）"]],
    [0.2, 0.5, 0.3], sizes=[10.5, 10, 10], rh=Inches(0.5), head_rh=Inches(0.44), first_col_color=GOLD)
rect(s, ML, Inches(5.65), CW, Inches(1.0), grad=[(0, "2A1E55", None), (100, "1E1542", None)], gang=120, line=GOLD, lw=1, radius=True)
rect(s, ML, Inches(5.65), Inches(0.07), Inches(1.0), fill=ROSE)
text(s, Inches(1.1), Inches(5.74), Inches(11.2), Inches(0.4), [("成本与人力（必列示）　", GOLD, True), ("企业服务暂不设专职，但属必做工作：按兼职/共享人力列示 约 15–30 万/年（社保+人力）。", INK, False)], size=11.5)
text(s, Inches(1.1), Inches(6.18), Inches(11.2), Inches(0.4), [("　　　　　　　　　　　", GOLD, False), ("第三方服务分成 按收入 30%–50%；系统/材料杂费 约 2–5 万/年。合理预估：首年企业服务收入(不含出海) ~30–60 万，净利 ~10–25 万。", MUT, False)], size=10.5, space=1.12)

# ---------- 25 章节09 规划拓展 ----------
nxt(); divider(IDX, "09", "规划建议与业务拓展", "PLANNING & EXPANSION", "规划设计 · 多元收入 · 资源导入 · 协同复制")

# ---------- 26 规划建议 ----------
nxt(); s = slide(); header(s, "09", "PLANNING & EXPANSION", "规划建议与业务拓展"); footer(s, IDX)
pl = [("规划设计", "配置 4–6 人设计条件；不建议做联合办公，突出产业办公与展贸", GOLD),
      ("多元化收入", "联合开展培训；利用小红书招募，增加经营性收入", ACC),
      ("资源导入", "引入复旦科技园进行客户导流，形成稳定客源", ACC2),
      ("空间利用", "顶层部分可举办漫展或各类展览（门票 + 赞助）", GREEN),
      ("协同复制", "合作模式可协同周边不超过两个项目，具备可复制性", ROSE),
      ("业务拓展", "以本项目为支点，进一步协助拓展新的合作项目", GOLD)]
for i, (t, d, c) in enumerate(pl):
    x = Emu(ML + (i % 3) * Inches(3.92)); y = Emu(Inches(1.85) + (i // 3) * Inches(2.35))
    card(s, x, y, Inches(3.72), Inches(2.15), title=t, body=d, accent=c, bsize=12.5)
    text(s, Emu(x + Inches(0.34)), Emu(y + Inches(1.5)), Inches(3), Inches(0.5), [("0%d" % (i + 1), c, True, FONT_EN)], size=30, font=FONT_EN)

# ---------- 27 章节09 -> 投决 ----------
nxt(); divider(IDX, "10", "投决建议与下一步", "DECISION & NEXT STEP", "三问决策 · 30 / 60 / 90 天 · 三年蓝图")

# ---------- 28 三问 + 路线 ----------
nxt(); s = slide(); header(s, "10", "DECISION & NEXT STEP", "三问决策 + 启动路线图"); footer(s, IDX)
q = [("Q1", "合作模式？", "轻资产运营（招商不对赌）：佣金 2–3 月 + 出海专项服务费 + 增值分成", GOLD),
     ("Q2", "报价档位？", "招商：佣金 2–3 月 + 活动 30 万 + 媒体 10 万 + 挂牌；出海：专项服务费 + 撮合佣金", ACC),
     ("Q3", "何时启动？", "与装修/设计交付同步，尽快签 MOU，先行政策梳理与挂牌筹备", ACC2)]
for i, (n, t, d, qc) in enumerate(q):
    x = Emu(ML + i * Inches(3.94))
    card(s, x, Inches(1.8), Inches(3.72), Inches(1.7), accent=qc)
    text(s, Emu(x + Inches(0.34)), Inches(2.0), Inches(3.2), Inches(0.4), [(n + "　", qc, True, FONT_EN), (t, INK, True)], size=15.5)
    text(s, Emu(x + Inches(0.34)), Inches(2.5), Inches(3.25), Inches(0.95), d, size=11, color=MUT, space=1.2)
road = [("30", "签 MOU + 招商手册与报价 + 政策汇编 + 首批挂牌对接"),
        ("60", "首场沙龙/路演 + OPC 社群与投流 + 首单招商 + 漫展筹备"),
        ("90", "正式挂牌 + 样板企业入驻 + 去化向首年 60% 推进")]
for i, (n, d) in enumerate(road):
    x = Emu(ML + i * Inches(3.94))
    card(s, x, Inches(3.7), Inches(3.72), Inches(1.7), accent=ACC)
    text(s, Emu(x + Inches(0.34)), Inches(3.9), Inches(2.6), Inches(0.6), [(n, ACC, True, FONT_EN), (" 天", MUT, False)], size=24)
    text(s, Emu(x + Inches(0.34)), Inches(4.55), Inches(3.25), Inches(0.85), d, size=11, color=MUT, space=1.18)
rect(s, ML, Inches(5.65), CW, Inches(0.95), grad=[(0, "322462", None), (100, "221848", None)], gang=0, line=GOLD, lw=1, radius=True)
text(s, Inches(1.1), Inches(5.65), Inches(11), Inches(0.95),
     [("三年蓝图　", GOLD, True), ("先签 6600㎡ 单项目专项（单项目单核算），三年跑通三年均 85% 去化后，协同周边 ≤2 项目并复制模式。", INK, False)],
     size=13.5, anchor=MSO_ANCHOR.MIDDLE)

# ---------- 29 封底 ----------
nxt(); s = slide()
oval(s, Inches(7.5), Inches(-2.4), Inches(8), Inches(8), grad=[(0, H_ACC, 20), (100, BG_A, 0)], gang=120)
oval(s, Inches(-2.5), Inches(3.0), Inches(6.5), Inches(6.5), grad=[(0, H_ACC2, 16), (100, BG_A, 0)], gang=60)
rect(s, 0, 0, Inches(0.16), SH, grad=[(0, H_GOLD, None), (100, H_ACC2, None)], gang=90)
text(s, ML, Inches(2.3), Inches(11), Inches(0.4), [("THANKS", GOLD, True, FONT_EN)], size=15, spacing=4)
rect(s, ML, Inches(2.85), Inches(0.9), Pt(3), fill=GOLD)
text(s, ML, Inches(3.1), Inches(11.5), Inches(1.8),
     [("让上海创智汇成为杨浦五角场\n", INK, True), ("AI + 数字内容", GOLD, True), ("的无界共创港与可复制样板", INK, True)], size=32, bold=True, space=1.2)
text(s, ML, Inches(5.1), Inches(11.5), Inches(0.5), "政策闭环  ·  三年合作  ·  协同复制  ·  内容自造流量", size=16, color=MUT, spacing=1)
rect(s, ML, Inches(6.4), CW, Pt(1), fill=LINE)
text(s, ML, Inches(6.55), Inches(11), Inches(0.4), [("谈判集中汇报 v3.0", GOLD, True), ("　期待与各方达成合作", SOFT, False)], size=12)

# 修正总页数
total = len(prs.slides._sldIdLst)
for r in _FOOT_RUNS:
    r.text = " / %02d" % total
ppt_path = os.path.join(HERE, "创智汇6600平合作方案.pptx")
prs.save(ppt_path)
print("PPT saved:", ppt_path, "slides:", total)
