# -*- coding: utf-8 -*-
"""创智汇 6600㎡ · 专题汇报版（回应中建关切）— 独立精简 PPT。
用法: python3 build_brief.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
FONT = "Microsoft YaHei"; FONT_EN = "Arial"
BG_A = "1A1140"; BG_B = "3A206E"
INK = RGBColor(0xF5, 0xF1, 0xFF); MUT = RGBColor(0xC6, 0xBA, 0xEA); SOFT = RGBColor(0x94, 0x86, 0xC4)
ACC = RGBColor(0x8B, 0x7B, 0xFF); ACC2 = RGBColor(0xC7, 0x7D, 0xFF)
GOLD = RGBColor(0xE9, 0xC2, 0x7C); GREEN = RGBColor(0x56, 0xD6, 0xB6); ROSE = RGBColor(0xF5, 0x8B, 0xB8)
LINE = RGBColor(0x4A, 0x3C, 0x7C); DARKTX = RGBColor(0x24, 0x14, 0x40)
H_ACC = "8B7BFF"; H_ACC2 = "C77DFF"; H_GOLD = "E9C27C"
SW, SH = Inches(13.333), Inches(7.5); ML = Inches(0.85); CW = Inches(11.63)
prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
BLANK = prs.slide_layouts[6]; _FOOT = []


def _grad(shape, stops, ang=90):
    spPr = shape._element.spPr
    for t in ('a:noFill', 'a:solidFill', 'a:gradFill', 'a:blipFill', 'a:pattFill', 'a:grpFill'):
        e = spPr.find(qn(t));  spPr.remove(e) if e is not None else None
    g = spPr.makeelement(qn('a:gradFill'), {}); lst = g.makeelement(qn('a:gsLst'), {})
    for pos, col, al in stops:
        gs = g.makeelement(qn('a:gs'), {'pos': str(int(pos * 1000))}); c = g.makeelement(qn('a:srgbClr'), {'val': col})
        if al is not None:
            c.append(g.makeelement(qn('a:alpha'), {'val': str(int(al * 1000))}))
        gs.append(c); lst.append(gs)
    g.append(lst); g.append(g.makeelement(qn('a:lin'), {'ang': str(int(ang * 60000)), 'scaled': '1'}))
    ln = spPr.find(qn('a:ln'));  (ln.addprevious(g) if ln is not None else spPr.append(g))


def slide(grad=True):
    s = prs.slides.add_slide(BLANK); bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.shadow.inherit = False; bg.line.fill.background()
    if grad:
        _grad(bg, [(0, BG_A, None), (55, BG_B, None), (100, BG_A, None)], 120)
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element); return s


def rect(s, x, y, w, h, fill=None, line=None, lw=1.0, radius=False, grad=None, gang=90):
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    b.shadow.inherit = False
    if grad is not None:
        _grad(b, grad, gang)
    elif fill is None:
        b.fill.background()
    else:
        b.fill.solid(); b.fill.fore_color.rgb = fill
    if line is None:
        b.line.fill.background()
    else:
        b.line.color.rgb = line; b.line.width = Pt(lw)
    return b


def oval(s, x, y, w, h, grad=None, gang=90):
    b = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h); b.shadow.inherit = False; b.line.fill.background()
    if grad is not None:
        _grad(b, grad, gang)
    return b


def text(s, x, y, w, h, runs, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=1.0, font=FONT, spacing=None):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    if isinstance(runs, str):
        runs = [(runs, color, bold)]
    p = tf.paragraphs[0]; p.alignment = align; p.line_spacing = space
    for seg in runs:
        t, c, b = seg[0], seg[1], seg[2]; f = seg[3] if len(seg) > 3 else font
        r = p.add_run(); r.text = t; r.font.size = Pt(size); r.font.bold = b; r.font.color.rgb = c; r.font.name = f
        if spacing is not None:
            r._r.get_or_add_rPr().set('spc', str(int(spacing * 100)))
    return tb


def bullets(s, x, y, w, h, items, size=14, color=MUT, gap=8, mark=ACC, lh=1.18):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True; tf.margin_left = Pt(2); tf.margin_top = Pt(2)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.space_after = Pt(gap); p.line_spacing = lh
        r0 = p.add_run(); r0.text = "▪  "; r0.font.size = Pt(size); r0.font.color.rgb = mark; r0.font.name = FONT
        r = p.add_run(); r.text = it; r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = FONT
    return tb


def header(s, sec, eyebrow, title):
    rect(s, ML, Inches(0.62), Inches(0.06), Inches(0.86), fill=GOLD)
    text(s, Emu(ML + Inches(0.22)), Inches(0.6), Inches(9), Inches(0.34), eyebrow, size=12.5, color=GOLD, bold=True, spacing=2, font=FONT_EN)
    text(s, Emu(ML + Inches(0.22)), Inches(0.92), Inches(10.4), Inches(0.66), title, size=26, color=INK, bold=True)
    text(s, Inches(11.3), Inches(0.5), Inches(1.55), Inches(0.9), [(sec, GOLD, True, FONT_EN)], size=40, align=PP_ALIGN.RIGHT, color=GOLD, font=FONT_EN)
    rect(s, Inches(11.55), Inches(1.32), Inches(1.28), Pt(2), fill=LINE)
    rect(s, ML, Inches(6.92), CW, Pt(1), fill=LINE)


def footer(s, idx):
    text(s, ML, Inches(7.0), Inches(9), Inches(0.32), [("上海创智汇 ", SOFT, False), ("CHUANGZHIHUI", SOFT, False, FONT_EN), ("  ·  6600㎡ 专题汇报版（回应中建关切）", SOFT, False)], size=9.5)
    rs = text(s, Inches(11.2), Inches(7.0), Inches(1.63), Inches(0.32), [("%02d" % idx, GOLD, True, FONT_EN), (" / 00", SOFT, False, FONT_EN)], size=10.5, align=PP_ALIGN.RIGHT)
    _FOOT.append(rs.text_frame.paragraphs[0].runs[1])


def card(s, x, y, w, h, title=None, body=None, items=None, accent=ACC, tsize=15.5, bsize=12.5, tcolor=INK):
    rect(s, x, y, w, h, grad=[(0, "2A1E55", None), (100, "1E1542", None)], gang=120, line=LINE, lw=1, radius=True)
    rect(s, x, y, Inches(0.07), h, fill=accent)
    if title:
        text(s, Emu(x + Inches(0.34)), Emu(y + Inches(0.2)), Emu(w - Inches(0.5)), Inches(0.5), title, size=tsize, color=tcolor, bold=True)
    if body:
        text(s, Emu(x + Inches(0.34)), Emu(y + Inches(0.66)), Emu(w - Inches(0.6)), Emu(h - Inches(0.78)), body, size=bsize, color=MUT, space=1.16)
    if items:
        bullets(s, Emu(x + Inches(0.34)), Emu(y + Inches(0.62)), Emu(w - Inches(0.6)), Emu(h - Inches(0.72)), items, size=bsize, mark=accent)


def table(s, x, y, w, rows, col_w, sizes=None, rh=Inches(0.52), head_rh=Inches(0.5), fcol=GOLD):
    ncol = len(rows[0]); sizes = sizes or [12.5] * ncol; cy = y
    for ri, row in enumerate(rows):
        ch = head_rh if ri == 0 else rh
        if ri == 0:
            rect(s, x, cy, w, ch, grad=[(0, "35276A", None), (100, "2A1E55", None)], gang=0)
        elif ri % 2 == 0:
            rect(s, x, cy, w, ch, fill=RGBColor(0x22, 0x18, 0x44))
        cx = x
        for ci, cell in enumerate(row):
            cwid = Emu(int(w * col_w[ci])); tb = s.shapes.add_textbox(cx, cy, cwid, ch); tf = tb.text_frame; tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_left = Pt(10); tf.margin_right = Pt(6); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
            p = tf.paragraphs[0]; p.line_spacing = 1.04; r = p.add_run(); r.text = str(cell)
            if ri == 0:
                r.font.size = Pt(12.5); r.font.bold = True; r.font.color.rgb = GOLD
            else:
                r.font.size = Pt(sizes[ci]); r.font.bold = (ci == 0); r.font.color.rgb = fcol if ci == 0 else MUT
            r.font.name = FONT; cx = Emu(cx + cwid)
        rect(s, x, Emu(cy + ch - Pt(0.75)), w, Pt(0.75), fill=LINE); cy = Emu(cy + ch)
    rect(s, x, y, Pt(2.5), Emu(cy - y), fill=GOLD); return cy


IDX = 0
def nxt():
    global IDX; IDX += 1; return IDX

# 1 封面
nxt(); s = slide()
oval(s, Inches(8.3), Inches(-2.0), Inches(7.5), Inches(7.5), grad=[(0, H_ACC, 22), (100, BG_A, 0)], gang=120)
oval(s, Inches(9.6), Inches(3.2), Inches(5.5), Inches(5.5), grad=[(0, H_ACC2, 18), (100, BG_A, 0)], gang=120)
rect(s, 0, 0, Inches(0.16), SH, grad=[(0, H_GOLD, None), (100, H_ACC2, None)], gang=90)
text(s, ML, Inches(1.05), Inches(11), Inches(0.4), [("BRIEFING · RESPONSE TO PARTNER", GOLD, True, FONT_EN)], size=13, spacing=3)
rect(s, ML, Inches(1.62), Inches(0.9), Pt(2.5), fill=GOLD)
text(s, ML, Inches(1.9), Inches(11.6), Inches(1.6), [("专题汇报 · 回应中建关切\n", INK, True), ("每个收费点：逻辑 · 价值 · 怎么做", GOLD, True)], size=36, bold=True, space=1.14)
text(s, ML, Inches(3.9), Inches(11.6), Inches(0.9),
     [("上海创智汇 · AI+数字内容无界共创港 · 6600㎡\n", MUT, False),
      ("回应：收费点具象化 / 活动 30 场 30 万 / 政策能否拿到现金 / 会客厅≥80 万 / 国际英才+运营", MUT, False)], size=15, space=1.5)
tags = ["每点：逻辑·价值·怎么做", "政策：能否拿到现金", "会客厅：≥80 万如何来", "出海：国际英才 + 运营"]
tx = ML
for t in tags:
    wd = Inches(0.42 + len(t) * 0.15); rect(s, tx, Inches(5.4), wd, Inches(0.52), line=LINE, lw=1.25, radius=True)
    text(s, tx, Inches(5.4), wd, Inches(0.52), t, size=11.5, color=INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE); tx = Emu(tx + wd + Inches(0.2))
rect(s, ML, Inches(6.4), CW, Pt(1), fill=LINE)
text(s, ML, Inches(6.52), Inches(11), Inches(0.4), [("汇报版 v1.0", GOLD, True), ("    用于与中建沟通；数据为合理测算，最终以协议与执行为准", SOFT, False)], size=11)

# 2 我们要回答的问题
nxt(); s = slide(); header(s, "00", "THE QUESTIONS", "我们要回答的问题（对方关切）"); footer(s, IDX)
qs = [("Q1", "每个收费点的背后逻辑与「可衡量价值」，以及具体怎么做？", GOLD),
      ("Q2", "活动 30 场/30 万 怎么做、什么形式、带来什么、是否带来收入？", ACC),
      ("Q3", "政策凭什么你们能申报？本项目可申哪些？能带来多少（能否拿到现金）？", ACC2),
      ("Q4", "国际会客厅收益逻辑与合作机制是什么？如何带来 ≥80 万？", GREEN)]
for i, (n, t, c) in enumerate(qs):
    y = Emu(Inches(1.95) + i * Inches(0.98))
    rect(s, ML, y, CW, Inches(0.82), grad=[(0, "2A1E55", None), (100, "1E1542", None)], gang=120, line=LINE, radius=True)
    rect(s, ML, y, Inches(0.07), Inches(0.82), fill=c)
    text(s, Emu(ML + Inches(0.3)), y, Inches(1.2), Inches(0.82), [(n, c, True, FONT_EN)], size=22, anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)
    text(s, Emu(ML + Inches(1.5)), y, Inches(9.9), Inches(0.82), t, size=14.5, color=INK, anchor=MSO_ANCHOR.MIDDLE)
rect(s, ML, Inches(5.95), CW, Inches(0.72), grad=[(0, "3A2878", None), (100, "2A1E55", None)], gang=0, line=GOLD, lw=1, radius=True)
text(s, Inches(1.1), Inches(5.95), Inches(11.2), Inches(0.72),
     [("补充　", GOLD, True), ("并把「国际英才（沪动营）+ 出海运营」补充进出海板块；本汇报逐条回应，给出逻辑、可衡量价值与落地打法。", INK, False)], size=12.5, anchor=MSO_ANCHOR.MIDDLE)

# 3 价值论证总览
nxt(); s = slide(); header(s, "Q1", "VALUE JUSTIFICATION", "收费点价值论证 · 逻辑 / 价值 / 怎么做"); footer(s, IDX)
table(s, ML, Inches(1.68), CW, [
    ["收费点", "背后逻辑（凭什么）", "给中建的可衡量价值", "怎么做"],
    ["招商佣金", "牌照+政策+社群自带客流，110:1 漏斗", "去化提速、租金回收提前", "牌照锚定+政策礼包+社群/活动导流"],
    ["活动 30 场/30 万", "高频活动引流+完成载体 KPI", "转化租客、补贴 10–20 万/年、媒体声量", "6 类活动打包，投流引流→线下转化"],
    ["媒体 10 万", "官媒背书 + 自媒体投流", "项目背书、招商曝光 ≥150 万", "官媒 5 万 + 自媒体/投流 5 万"],
    ["政策申报", "挂牌科企服务中心+科企联+复旦系资质", "企业补贴到账、平台服务费/分成", "见 Q3 政策页"],
    ["企业服务", "联合公司承接注册/申报/知产/财税", "企业黏性、第二收入曲线", "联合公司经营分成"],
    ["出海撮合", "领事/大湾区/高校/国际英才资源", "海外订单、租金反哺、合资收益", "全链条 + 领事路径"],
    ["国际会客厅", "领事资源合作方 + 高端背书", "外资对接、高端招商、≥80 万收益", "见 Q4 会客厅页"]],
    [0.14, 0.29, 0.31, 0.26], sizes=[10.5, 9.5, 9.5, 9.5], rh=Inches(0.52), head_rh=Inches(0.44))
text(s, ML, Inches(6.55), Inches(11.6), Inches(0.3), "＊每个收费点＝背后逻辑（我们凭什么）+ 给中建可衡量价值 + 具体怎么做。", size=9.5, color=SOFT)

# 4 活动
nxt(); s = slide(); header(s, "Q2", "EVENT VALUE", "活动 30 场/30 万 · 怎么做 · 效果 · 是否带收入"); footer(s, IDX)
card(s, ML, Inches(1.75), Inches(3.78), Inches(3.05), "怎么做（形式）", accent=ACC, bsize=11.5, items=[
    "AI 培训 / 黑客松、政策沙龙 / 路演", "行业对接 / 撮合 / 游学", "IP 潮玩市集、漫展 / 大型活动", "领事专题 / 买家团 / 出海推介", "线上投流引流 → 线下办会 → 转化"])
card(s, Inches(4.79), Inches(1.75), Inches(3.78), Inches(3.05), "带来什么效果（可衡量）", accent=GOLD, bsize=11.5, items=[
    "每场 ≥30 人，年触达 900+ 精准客户", "110:1 招商漏斗，转化租客/企业", "完成载体 KPI → 补贴 10–20 万/年", "媒体曝光 ≥150 万，提升声量", "为出海/撮合/会籍持续输送客源"])
card(s, Inches(8.78), Inches(1.75), Inches(3.7), Inches(3.05), "是否带来收入", accent=GREEN, bsize=11.5, items=[
    "漫展：门票 + 赞助（经营性收入）", "培训 / 训练营：按人/按场收费", "赞助商 / 展位 / 冠名", "活动转化企业服务与撮合收入", "→ 部分自造收入冲抵活动成本"])
rect(s, ML, Inches(5.0), CW, Inches(1.55), grad=[(0, "2A1E55", None), (100, "1E1542", None)], gang=120, line=GOLD, lw=1, radius=True)
rect(s, ML, Inches(5.0), Inches(0.07), Inches(1.55), fill=GOLD)
text(s, Inches(1.1), Inches(5.12), Inches(11), Inches(0.4), "给中建的价值 & 费用说明", size=13.5, color=GOLD, bold=True)
bullets(s, Inches(1.1), Inches(5.55), Inches(5.6), Inches(0.9), ["价值：去化提速 + 补贴达标 + 品牌声量 + 企业转化", "30 场为满足招商/载体/补贴各项要求的基线，可按招人方式调整"], size=11, gap=4, mark=GOLD)
bullets(s, Inches(6.9), Inches(5.55), Inches(5.5), Inches(0.9), ["30 万为固定活动执行费（策划/执行/导师/物料）", "门票/赞助/培训等经营性收入单独结算，可冲抵成本"], size=11, gap=4, mark=GOLD)

# 5 政策能力+清单
nxt(); s = slide(); header(s, "Q3", "POLICY CAPABILITY", "政策申报 · 凭什么我们能申 + 本项目可申清单"); footer(s, IDX)
card(s, ML, Inches(1.7), Inches(5.72), Inches(2.5), "凭什么我们能申报（能力/资质）", accent=GOLD, bsize=10.5, items=[
    "挂牌科企服务中心 + 依托科企联，政务快通道", "拟入驻一网通办创新券平台，作服务机构核销",
    "复旦系（技术转移/珠海院/成果转化中心）+ 高校", "专业申报团队 + 过往经验（高企/专精特新）", "载体认定路径：载体/成果转化平台/OPC 社区"])
card(s, Inches(6.76), Inches(1.7), Inches(5.72), Inches(2.5), "本项目可申报清单（具体）", accent=ACC, bsize=10.5, items=[
    "企业侧：高企 20 万、专精特新 10/30 万、科小备案", "企业侧：AI 房租补贴 ≤100 万/年、算力 ≤50%、创新券 ≤30 万/年",
    "企业侧：微短剧专项（智能体/优秀剧/厂牌/出海）", "载体侧：成果转化平台运营 ≤100 万/年、服务业引导资金 ≤300 万", "载体侧：活动/载体补贴（满年 10–20 万）、载体认定"])
rect(s, ML, Inches(4.38), CW, Inches(2.05), grad=[(0, "2A1E55", None), (100, "1E1542", None)], gang=120, line=GOLD, lw=1, radius=True)
rect(s, ML, Inches(4.38), Inches(0.07), Inches(2.05), fill=GREEN)
text(s, Inches(1.1), Inches(4.5), Inches(11.2), Inches(0.4), "政策 → 收入怎么算（平台可实现现金）", size=13, color=GREEN, bold=True)
bullets(s, Inches(1.1), Inches(4.95), Inches(5.6), Inches(1.4), ["申报服务费：科小 0.3–2 万/家；高企/专精特新 2–8 万/家", "成功分成：补贴到账额 5%–15%（以兑现为准）", "载体运营补贴：认定+满年约 10–20 万/年（直接到平台）"], size=10.5, gap=5, mark=GREEN)
bullets(s, Inches(6.9), Inches(4.95), Inches(5.5), Inches(1.4), ["前提：先取得相应认定（载体/服务机构），摸排已列", "口径：补贴主体到企业/载体，平台按服务费+分成取酬", "先陪 3–5 家样板跑通，形成可复制申报流水"], size=10.5, gap=5, mark=GREEN)

# 6 政策→现金 + 变现前置
nxt(); s = slide(); header(s, "Q3", "POLICY → CASH", "政策能否拿到现金 · 变现前置条件"); footer(s, IDX)
table(s, ML, Inches(1.7), CW, [
    ["政策", "金额（含单位）", "平台可实现现金", "现金确定性"],
    ["载体运营/活动补贴（认定后）", "≤100 万/年；满年约 10–20 万/年", "直接到平台 10–20 万/年", "现金·较确定（需认定+满年）"],
    ["高企 / 专精特新 / 科小", "20 万 / 10·30 万 / 资格（企业）", "服务费 0.3–8 万/家", "现金·确定（服务费）"],
    ["AI 房租补贴 / 政府补贴项目", "≤100 万/年 / 视项目（企业）", "成功费 到账 5%–15%", "现金·以兑现为准"],
    ["创新券", "企业≤30 万/年（抵付服务）", "服务机构核销分成", "现金·需入驻平台"],
    ["算力补贴 + 云折扣", "最高 50%（企业）", "招商抓手（非平台现金）", "非现金（招商用）"]],
    [0.28, 0.28, 0.24, 0.2], sizes=[10.5, 10, 10, 9.5], rh=Inches(0.5), head_rh=Inches(0.44))
rect(s, ML, Inches(4.9), CW, Inches(1.7), grad=[(0, "3A2878", None), (100, "2A1E55", None)], gang=0, line=GOLD, lw=1.25, radius=True)
text(s, Inches(1.1), Inches(5.0), Inches(11.2), Inches(0.4), "政策变现前置开关（建议与中建对齐）", size=13.5, color=GOLD, bold=True)
bullets(s, Inches(1.1), Inches(5.45), Inches(5.6), Inches(1.1), ["① 项目/平台取得「经认定载体 / 成果转化服务平台」→ 载体补贴到平台", "② 我方入驻「一网通办创新券平台」作服务机构 → 核销分成"], size=11, gap=5, mark=GOLD)
bullets(s, Inches(6.9), Inches(5.45), Inches(5.5), Inches(1.1), ["③ 首批 3–5 家样板企业申报跑通 → 形成收入流水", "计入预算的政策收入＝平台可实现现金（服务费+分成+载体补贴）"], size=11, gap=5, mark=GOLD)

# 7 会客厅
nxt(); s = slide(); header(s, "Q4", "SALON LOGIC", "国际会客厅 · 收益逻辑与合作机制（≥80 万）"); footer(s, IDX)
rect(s, ML, Inches(1.68), CW, Inches(1.0), grad=[(0, "3A2878", None), (100, "2A1E55", None)], gang=0, line=GOLD, lw=1.25, radius=True)
text(s, Inches(1.1), Inches(1.68), Inches(11.2), Inches(1.0),
     [("合作机制　", GOLD, True), ("三方协同：平台（运营策划+活动执行+出海撮合）× 领事资源合作方（领事网络+空间）× 复旦政策研究中心/科企联（背书+获客）；以「国别/北欧会客厅」为载体，按 冠名+会籍+活动+出海对接 收费，收益按项目分润。", INK, False)],
     size=12, anchor=MSO_ANCHOR.MIDDLE, space=1.16)
text(s, ML, Inches(2.85), Inches(11), Inches(0.35), "≥80 万收益如何构成（保守单一情景）", size=13.5, color=GOLD, bold=True)
table(s, ML, Inches(3.3), CW, [
    ["收益项", "口径", "年收益（保守）"],
    ["国家/联合冠名", "1 个 × 30–80 万/年", "约 40 万"],
    ["企业常驻会籍", "3–5 家 × 10 万/家/年", "约 30–50 万"],
    ["单场活动（领事出席）", "6–8 场 × 3–8 万/场", "约 20–40 万"],
    ["出海对接服务", "单国服务包 20–50 万/项目 × 1–2 单", "约 20–50 万"],
    ["合计（毛收益）", "以上累计", "约 110–180 万 → 分润/成本后平台 ≥80 万"]],
    [0.26, 0.46, 0.28], sizes=[11, 10.5, 11], rh=Inches(0.5), head_rh=Inches(0.44))
text(s, ML, Inches(6.5), Inches(11.6), Inches(0.3), "＊给中建价值：以「国家为单位」的招商与外资对接入口、高端背书；收益为保守测算，最终以合作方协议为准。", size=9.5, color=SOFT)

# 8 国际英才 + 运营
nxt(); s = slide(); header(s, "+", "GLOBAL TALENT & OPS", "出海板块补充 · 国际英才 + 出海运营"); footer(s, IDX)
rect(s, ML, Inches(1.7), Inches(7.75), Inches(4.85), grad=[(0, "2A1E55", None), (100, "1E1542", None)], gang=120, line=GOLD, lw=1.25, radius=True)
rect(s, ML, Inches(1.7), Inches(0.07), Inches(4.85), fill=GOLD)
text(s, Inches(1.08), Inches(1.84), Inches(7.4), Inches(0.4), "国际英才 · 沪动营（留学生丝路英才计划）", size=15, color=GOLD, bold=True)
text(s, Inches(1.08), Inches(2.3), Inches(7.4), Inches(0.5), [("定位　", ACC, True), ("「沪联世界·智创未来」——财大/复旦/同济三校联合，服务逾万名在校留学生，培养懂中国·通国际的一带一路桥梁人才。", MUT, False)], size=11, space=1.15)
bullets(s, Inches(1.08), Inches(3.05), Inches(7.5), Inches(2.2), [
    "三大体系：思想引领（丝路青年领袖峰会）→ 能力提升（1+N 双导师·跨境电商实训）→ 实践转化（跨境贸易服务中心·全球供应链）",
    "出海用途：双语双市场人才、TikTok 直播选品、海外品牌孵化、连接一带一路国家母国市场",
    "百千万工程（3 年）：孵化 100 跨境团队 / 促成 1000 人就业 / 培养 10000 复合人才；目标 10 亿跨境贸易、30 海外品牌",
    "收入落点：训练营/创投大赛/跨境撮合佣金/企业结对 + 留学生政策项目"], size=10.5, gap=6, mark=ACC)
card(s, Inches(8.95), Inches(1.7), Inches(3.53), Inches(4.85), "出海运营（运营）", accent=ACC2, bsize=11, items=[
    "出海代运营：领英/谷歌/TikTok/短剧", "海外账号搭建与投放代运营", "出海活动运营 + 买家团/展销运营", "落地陪跑：连领事→部长→洽谈→落地", "团队：初期兼职 → 成熟全职约 3 人", "收费：出海专项服务费 + 代运营 + 撮合佣金"])

# 9 结论与下一步
nxt(); s = slide(); header(s, "END", "CONCLUSION & NEXT", "结论与下一步（与中建对齐）"); footer(s, IDX)
card(s, ML, Inches(1.75), Inches(5.72), Inches(2.6), "结论：每点都讲得清", accent=GOLD, bsize=12, items=[
    "每个收费点＝逻辑 + 可衡量价值 + 落地打法", "活动：形式/效果/收入清晰，30 万为固定执行费", "政策：平台现金＝服务费+成功分成+载体补贴", "会客厅：机制清晰，≥80 万有构成拆解", "出海：国际英才 + 运营 双支撑，非空谈"])
card(s, Inches(6.76), Inches(1.75), Inches(5.72), Inches(2.6), "下一步：先开前置开关", accent=GREEN, bsize=12, items=[
    "① 推进载体/成果转化平台认定", "② 我方入驻创新券平台作服务机构", "③ 锁定首批 3–5 家样板企业", "④ 确认会客厅合作方与分润机制", "⑤ 锁定出海国家/资源与运营团队"])
rect(s, ML, Inches(4.55), CW, Inches(2.0), grad=[(0, "3A2878", None), (100, "2A1E55", None)], gang=0, line=GOLD, lw=1.25, radius=True)
text(s, Inches(1.1), Inches(4.68), Inches(11.2), Inches(0.4), "一句话对齐", size=14, color=GOLD, bold=True)
text(s, Inches(1.1), Inches(5.12), Inches(11.2), Inches(1.3),
     [("方案不是「只有报价」——", INK, True),
      ("每个收费点背后都有明确逻辑、可衡量价值与落地打法；政策收入以「能拿到现金」口径计入（先认定、再兑现）；国际会客厅有清晰机制与 ≥80 万收益构成；出海以国际英才与运营团队落到实处。建议本周对齐三个认定/合作前置开关，即可进入预算收口与签约。", MUT, False)],
     size=12.5, space=1.3)

total = len(prs.slides._sldIdLst)
for r in _FOOT:
    r.text = " / %02d" % total
p = os.path.join(HERE, "创智汇6600平-汇报版.pptx")
prs.save(p); print("saved", p, "slides", total)
