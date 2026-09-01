# -*- coding: utf-8 -*-
"""创智汇 · 30场活动专项交付（修订版 v2）
仅活动专项；不做租金租赁对赌。
口径：
- 每场人次 ≤30；负责人占比≤30%
- 关联度弱化表述，删除核验表交付
- 8月改为项目推广日（中介/五大行/投研/政府等）
- 潮玩/毛绒等改为 ChinaJoy 主题内容
- 活动衔接 WAIC + ChinaJoy
- 出海活动另议另计价（不写「慎重」字样）
- 领事到访/挂牌另计价（在三部分费用之外）
- 不承诺必带外资企业
- 同浦汇带客，园区负责销售；群内协同逼单
- 付款：签约50% + 四季各10% + 年终10%；取消方案二
主体：云创基地（国家级孵化器）纳入运营范围；
学术支持：复旦住房政策研究中心；
活动及运营载体支持：科企联 + 科企服务中心。
"""
import os, shutil
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

HERE = os.path.dirname(os.path.abspath(__file__))
ART = "/opt/cursor/artifacts"
PHOTO = os.path.join(HERE, "cases", "photos")
FONT = "Microsoft YaHei"

INK = RGBColor(0x1A, 0x1A, 0x2E)
MUT = RGBColor(0x4A, 0x4A, 0x5A)
SOFT = RGBColor(0x8A, 0x8A, 0x9A)
BLUE = RGBColor(0x1E, 0x4D, 0x8C)
GOLD = RGBColor(0xC4, 0x9A, 0x3C)
LINE = RGBColor(0xD8, 0xDE, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF7, 0xF8, 0xFA)
ACC = RGBColor(0xE8, 0xF0, 0xFA)
GREEN = RGBColor(0x1F, 0x7A, 0x5C)
ROSE = RGBColor(0xA6, 0x3D, 0x40)

SW, SH = Inches(13.333), Inches(7.5)
ML, CW = Inches(0.55), Inches(12.23)

PARTIES = {
    "ops": "上海市云计算创新基地（国家级孵化器）——纳入主体运营范围",
    "academic": "复旦大学住房政策研究中心——学术支持单位",
    "support": "上海市杨浦区科技企业联合会、上海市杨浦区科技企业服务中心——活动及运营载体支持单位",
    "exec": "同浦汇——活动策划执行与带客",
}

# 30场：人次≤30；负责人≤30%；ChinaJoy主题替换潮玩毛绒；8月项目推广日；出海另议
EVENTS = [
    # —— 8月：项目推广日（中介/五大行/投研/政府）——
    dict(code="P1", cat="P", name="创智汇项目推广日①·中介与渠道专场", month="2026年8月上旬", fmt="项目推广日",
         people="≤30人", dm="≤30%", floor="5F+3F", rhythm="T-14物料→T-7邀约→T-1彩排→D日→建群回访",
         lease="中介渠道带看与项目推介",
         agenda="项目介绍→空间导览→政策要点→一对一洽谈→建群",
         guests="房产/产业中介、渠道机构、园区招商同事", value="渠道点火；待租单元曝光",
         deliver="方案+议程+签到+意向表+建群+回访摘要+复盘"),
    dict(code="P2", cat="P", name="创智汇项目推广日②·五大行与金融机构专场", month="2026年8月中旬", fmt="项目推广日",
         people="≤30人", dm="≤30%", floor="5F沙龙", rhythm="同上",
         lease="金融资源链接与企业客户导入",
         agenda="项目路演→金融产品对接→企业需求座谈→看场→建群",
         guests="五大行及合作金融机构对公/科创条线", value="金融背书；客户转介",
         deliver="同上"),
    dict(code="P3", cat="P", name="创智汇项目推广日③·投研机构专场", month="2026年8月中旬", fmt="项目推广日",
         people="≤30人", dm="≤30%", floor="5F圆桌", rhythm="同上",
         lease="投研视角放大项目影响力",
         agenda="项目定位→产业赛道研判→圆桌→看场→建群",
         guests="券商研究所、产业研究院、智库分析师", value="专业传播；精准客群",
         deliver="同上"),
    dict(code="P4", cat="P", name="创智汇项目推广日④·政府部门与载体协同专场", month="2026年8月下旬", fmt="项目推广日",
         people="≤30人", dm="≤30%", floor="3F+5F", rhythm="同上",
         lease="政策与载体协同落地",
         agenda="载体介绍→政策协同→云创基地资源→看场→建群",
         guests="区相关部门、载体代表、云创基地、服中心", value="政策协同；主体背书",
         deliver="同上"),
    dict(code="E1", cat="E", name="WAIC成果承接·创智汇开放交流日", month="2026年8月下旬", fmt="开放交流日",
         people="≤30人", dm="≤30%", floor="3F+5F", rhythm="同上",
         lease="衔接WAIC议题与园区看场",
         agenda="WAIC议题速递→园区承接→双展台→洽谈→建群",
         guests="AI/内容企业、媒体、云创基地代表", value="WAIC衔接；看场转化",
         deliver="同上+通稿"),
    # —— 常规训练/政策/对接 ——
    dict(code="A2", cat="A", name="Agent智能体搭建一日营", month="2026年9月上旬", fmt="实训营",
         people="≤30人", dm="≤30%", floor="3F OPC", rhythm="同上",
         lease="Agent初创落位", agenda="签到→工作流→实操→Demo→看场",
         guests="Agent工程师、初创团队", value="入孵转化", deliver="同上"),
    dict(code="C1", cat="C", name="ChinaJoy主题·数字娱乐与IP授权对接会", month="2026年9月中旬", fmt="对接会",
         people="≤30人", dm="≤30%", floor="5F展示中心", rhythm="同上",
         lease="ChinaJoy赛道企业/IP方入驻与展位",
         agenda="ChinaJoy赛道解读→IP授权→一对一撮合→看场→建群",
         guests="游戏/动漫/潮玩IP方、渠道、内容公司", value="衔接ChinaJoy；展位/办公转化",
         deliver="同上"),
    dict(code="F1", cat="F", name="出海专题活动（另议·另计价）", month="档期另议", fmt="专题（另议）",
         people="≤30人", dm="≤30%", floor="另议", rhythm="立项后单独排期",
         lease="出海需求企业服务（不在年度活动包费用内）",
         agenda="另案确定", guests="另案确定",
         value="出海场次价格与方案另议，不计入本册活动服务费",
         deliver="另案合同与交付清单"),
    dict(code="A3", cat="A", name="多模态智能体工作坊", month="2026年10月上旬", fmt="工作坊",
         people="≤30人", dm="≤30%", floor="3F培训", rhythm="同上",
         lease="多模态团队入驻", agenda="案例→实操→共创→评审→看场",
         guests="多模态应用团队、产品负责人", value="技术团队入驻", deliver="同上"),
    dict(code="E4", cat="E", name="ChinaJoy主题日·数字娱乐生态交流", month="2026年10月中旬", fmt="主题日",
         people="≤30人", dm="≤30%", floor="5F", rhythm="同上",
         lease="ChinaJoy生态企业集中看场与洽谈",
         agenda="开幕→赛道分享→展洽→看场→建群",
         guests="数字娱乐企业、内容工作室、渠道", value="衔接ChinaJoy；集中洽谈",
         deliver="同上"),
    dict(code="B2", cat="B", name="AI治理与可信智能体沙龙", month="2026年10月下旬", fmt="沙龙",
         people="≤30人", dm="≤30%", floor="3F", rhythm="同上",
         lease="合规型企业信任导入", agenda="治理框架→合规清单→案例→问答→对接",
         guests="合规顾问、企业法务/负责人", value="信任入驻", deliver="同上"),
    dict(code="A4", cat="A", name="火山引擎×算力Infra实务营", month="2026年11月上旬", fmt="厂商联训",
         people="≤30人", dm="≤30%", floor="3F+云创基地", rhythm="同上",
         lease="高算力企业定向看场", agenda="政策包→代金券→案例→诊断→看场",
         guests="火山引擎、云创基地、企业技术负责人", value="云创协同；算力企业", deliver="同上"),
    dict(code="B4", cat="B", name="创新券·算力券·模型券实务沙龙", month="2026年11月中旬", fmt="实务沙龙",
         people="≤30人", dm="≤30%", floor="3F", rhythm="同上",
         lease="用券企业聚集", agenda="三券规则→核销→案例→开户→入驻洽谈",
         guests="券平台、云厂商、企业负责人", value="券务转化", deliver="同上"),
    dict(code="L1", cat="L", name="领事到访接待（另计价）", month="档期另议", fmt="外事接待（另计价）",
         people="另议", dm="—", floor="会客厅/5F", rhythm="外事流程单独排期",
         lease="领事到访接待（不在三部分活动费用内）",
         agenda="另案确定", guests="领事及相关方（另案）",
         value="领事到访与挂牌活动分属不同性质，均另计价",
         deliver="另案报价与执行方案"),
    dict(code="A5", cat="A", name="具身智能空间交互体验日", month="2026年12月上旬", fmt="体验日",
         people="≤30人", dm="≤30%", floor="5F+3F", rhythm="同上",
         lease="具身/机器人团队看场", agenda="演示→讲解→场景→踩点→洽谈",
         guests="具身团队、高校实验室", value="具身看场", deliver="同上"),
    dict(code="F4", cat="F", name="创智汇AI年度Demo Day", month="2026年12月中旬", fmt="路演日",
         people="≤30人", dm="≤30%", floor="5F主展", rhythm="同上",
         lease="集中签约洽谈与媒体背书", agenda="开幕→精选路演→颁奖→洽谈",
         guests="投资人、链主、媒体、精选项目", value="集中签约", deliver="同上+签约台账"),
    dict(code="D1", cat="D", name="AIGC微短剧制片特训", month="2027年1月上旬", fmt="特训营",
         people="≤30人", dm="≤30%", floor="3F+5F", rhythm="同上",
         lease="厂牌/工作室入驻", agenda="政策→制片→脚本→路演→看场",
         guests="导演、厂牌、平台方", value="内容入驻", deliver="同上"),
    dict(code="A6", cat="A", name="AI营销Agent实战营", month="2027年1月中旬", fmt="实战营",
         people="≤30人", dm="≤30%", floor="3F", rhythm="同上",
         lease="营销科技公司看场", agenda="策略→Agent→素材→复盘→转化",
         guests="投放操盘手、产品经理", value="营销科技入驻", deliver="同上"),
    dict(code="B1", cat="B", name="YOUNG立方×智能伙伴政策沙龙", month="2027年2月中旬", fmt="政策沙龙",
         people="≤30人", dm="≤30%", floor="5F沙龙", rhythm="同上",
         lease="内容/AI企业导入", agenda="政策→礼包→画像→诊断→看场",
         guests="政策宣讲、服中心、企业负责人", value="政策获客", deliver="同上"),
    dict(code="B3", cat="B", name="高企认定冲刺（AI企业专场）", month="2027年2月下旬", fmt="辅导会",
         people="≤30人", dm="≤30%", floor="3F", rhythm="同上",
         lease="待认定企业带政策入驻", agenda="条件→材料→初筛→套餐→激励",
         guests="辅导顾问、待认定企业负责人", value="带政策入驻", deliver="同上"),
    dict(code="A7", cat="A", name="OPC超级个体黑客松（春）", month="2027年3月上旬", fmt="黑客松",
         people="≤30人", dm="≤30%", floor="3F", rhythm="同上",
         lease="获奖团队优先谈单元", agenda="开题→开发→路演→礼包→看场",
         guests="评委、投资人、Builders", value="获奖谈单元", deliver="同上"),
    dict(code="C2", cat="C", name="高校成果转化·AI for Science日", month="2027年3月下旬", fmt="对接日",
         people="≤30人", dm="≤30%", floor="3F+沙龙", rhythm="同上",
         lease="成果公司/实验室落户", agenda="成果路演→场景→承接→看单元→洽谈",
         guests="复旦/同济技转、教授团队、住房政策研究中心学者（学术支持）", value="成果落户", deliver="同上"),
    dict(code="F3", cat="F", name="通往AGI季度圆桌", month="2027年3月下旬", fmt="闭门圆桌",
         people="≤30人", dm="≤30%", floor="5F沙龙", rhythm="同上",
         lease="研究型/模型团队", agenda="议题→圆桌→围炉→问答→看场",
         guests="学者、模型企业、投资人", value="研究型客户", deliver="同上"),
    dict(code="D2", cat="D", name="ChinaJoy主题·创作者内容首发①", month="2027年4月中旬", fmt="发布会",
         people="≤30人", dm="≤30%", floor="5F", rhythm="同上",
         lease="数字娱乐/内容品牌问询", agenda="揭幕→发布→快闪→专访→招商通道",
         guests="创作者、媒体、渠道、ChinaJoy相关内容方", value="衔接ChinaJoy内容侧", deliver="同上"),
    dict(code="C3", cat="C", name="ChinaJoy主题·游戏周边与潮玩供应链对接", month="2027年4月下旬", fmt="对接会",
         people="≤30人", dm="≤30%", floor="5F", rhythm="同上",
         lease="ChinaJoy供应链/周边企业展位落位",
         agenda="赛道介绍→供需对接→展位参观→报价→建群",
         guests="游戏周边、潮玩供应链、渠道采购", value="ChinaJoy供应链去化", deliver="同上"),
    dict(code="C4", cat="C", name="专精特新·AI应用培育路演", month="2027年5月中旬", fmt="路演",
         people="≤30人", dm="≤30%", floor="5F", rhythm="同上",
         lease="成长型AI补位", agenda="政策→路演→点评→金融→看场",
         guests="顾问、银行、基金、企业", value="成长型企业", deliver="同上"),
    dict(code="D3", cat="D", name="ChinaJoy主题·数字娱乐市集体验日", month="2027年5月下旬", fmt="体验日",
         people="≤30人", dm="≤30%", floor="5F", rhythm="同上",
         lease="优质摊主/工作室升级固定展位",
         agenda="布展→体验→交流→转正洽谈→建群",
         guests="内容工作室、周边品牌、达人", value="ChinaJoy氛围转化", deliver="同上"),
    dict(code="C5", cat="C", name="ChinaJoy主题·IP联名与衍生品撮合会", month="2027年6月中旬", fmt="撮合会",
         people="≤30人", dm="≤30%", floor="5F", rhythm="同上",
         lease="IP联名/衍生品团队入驻",
         agenda="IP路演→联名模式→一对一→MOU→看场",
         guests="IP方、衍生品商、渠道", value="ChinaJoy IP商业化落地", deliver="同上"),
    dict(code="D4", cat="D", name="ChinaJoy主题·沉浸式数字娱乐联展", month="2027年6月下旬", fmt="联展（短展期）",
         people="≤30人/场次接待", dm="≤30%", floor="5F", rhythm="展期预约制",
         lease="闭幕洽谈集中转化", agenda="布展→预约观展→交流→闭幕洽谈",
         guests="联合内容方、渠道、品牌", value="ChinaJoy视觉落地", deliver="同上"),
    dict(code="D5", cat="D", name="创作者首发②·衔接WAIC2027与ChinaJoy", month="2027年7月中旬", fmt="发布会",
         people="≤30人", dm="≤30%", floor="5F", rhythm="同上",
         lease="旺季补位，双大会预热", agenda="发布→渠道→看场→年框→双会预热",
         guests="IP方、渠道、媒体", value="同步衔接WAIC与ChinaJoy", deliver="同上"),
]

assert len(EVENTS) == 30

CASES = [
    {
        "name": "杨「数」浦数字沙龙：AI如何重塑企业DNA",
        "meta": "2025-06-30 · 美团上海综合指挥中心 · 沙龙类",
        "orgs": "杨浦区委网信办；赛博院等",
        "companies": "美团；大众点评大模型团队等",
        "agenda": ["参观数字化场景", "主题分享", "案例落地", "政策解读", "互动答疑"],
        "photos": ["case_c987bd91.jpg", "case_59157ea5.jpg", "case_6e51af9e.jpg"],
        "map": "对应：B类政策沙龙、A类训练营",
    },
    {
        "name": "「融见科创」人工智能专场路演暨投融资对接会",
        "meta": "2025-10 · 杨浦 · 路演对接类",
        "orgs": "杨浦科创促进会、邮储银行等",
        "companies": "复楚智能、卡房信息、一造科技、中科趋势、万笔千墨等",
        "agenda": ["开场致辞", "主办致辞", "趋势分享", "项目路演对接", "点评跟进"],
        "photos": ["case_db1d5cac.jpg", "case_65931516.jpg", "case_cb0a3812.jpg"],
        "map": "对应：F4 Demo Day、项目推广日、路演类",
    },
]


def rect(s, x, y, w, h, fill=None, line=None, lw=0.8):
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    b.shadow.inherit = False
    if fill is None: b.fill.background()
    else: b.fill.solid(); b.fill.fore_color.rgb = fill
    if line is None: b.line.fill.background()
    else: b.line.color.rgb = line; b.line.width = Pt(lw)
    return b


def text(s, x, y, w, h, content, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = content; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = FONT
    return tb


def bullets(s, x, y, w, h, items, size=12, color=MUT, mark=BLUE):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4); p.line_spacing = 1.15
        r0 = p.add_run(); r0.text = "▪ "; r0.font.size = Pt(size); r0.font.color.rgb = mark; r0.font.name = FONT
        r = p.add_run(); r.text = it; r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = FONT
    return tb


def slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, SW, SH, fill=WHITE)
    rect(s, 0, 0, SW, Inches(0.06), fill=BLUE)
    return s


def header(s, sec, title, sub=None):
    text(s, ML, Inches(0.22), Inches(12), Inches(0.28), sec, size=12, color=BLUE, bold=True)
    text(s, ML, Inches(0.48), Inches(12), Inches(0.42), title, size=22, color=INK, bold=True)
    if sub:
        text(s, ML, Inches(0.92), Inches(12), Inches(0.28), sub, size=12, color=MUT)


def footer(s, page, total=20):
    text(s, ML, Inches(7.15), Inches(9.2), Inches(0.25),
         "上海创智汇 × 同浦汇 · 活动专项 · 衔接 WAIC & ChinaJoy", size=10, color=SOFT)
    text(s, Inches(11.3), Inches(7.15), Inches(1.5), Inches(0.25), f"{page}/{total}", size=10, color=SOFT, align=PP_ALIGN.RIGHT)


def card(s, x, y, w, h, title, items, accent=BLUE):
    rect(s, x, y, w, h, fill=BG, line=LINE)
    rect(s, x, y, Inches(0.08), h, fill=accent)
    text(s, Emu(x + Inches(0.22)), Emu(y + Inches(0.12)), Emu(w - Inches(0.35)), Inches(0.32), title, size=14, color=INK, bold=True)
    bullets(s, Emu(x + Inches(0.22)), Emu(y + Inches(0.48)), Emu(w - Inches(0.4)), Emu(h - Inches(0.55)), items, size=12, mark=accent)


def table(s, y0, rows, col_w, sizes=None, rh=0.32):
    yy = y0
    if sizes is None: sizes = [10] * len(rows[0])
    for ri, row in enumerate(rows):
        xx = ML; h = Inches(0.38 if ri == 0 else rh)
        for ci, val in enumerate(row):
            cw = Emu(int(CW * col_w[ci]))
            bg = BLUE if ri == 0 else (ACC if ri % 2 == 0 else WHITE)
            fg = WHITE if ri == 0 else INK
            rect(s, xx, yy, cw, h, fill=bg, line=LINE, lw=0.5)
            text(s, Emu(xx + Inches(0.04)), Emu(yy + Inches(0.04)), Emu(cw - Inches(0.06)), Emu(h - Inches(0.06)),
                 str(val), size=sizes[ci], color=fg, bold=(ri == 0), anchor=MSO_ANCHOR.MIDDLE)
            xx = Emu(xx + cw)
        yy = Emu(yy + h)
    return yy


def add_pic(s, path, x, y, w, h):
    if os.path.exists(path):
        s.shapes.add_picture(path, x, y, width=w, height=h)
    else:
        rect(s, x, y, w, h, fill=ACC, line=LINE)


def build():
    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
    TOTAL = 20

    # 1 cover
    s = slide(prs)
    text(s, ML, Inches(1.5), Inches(12), Inches(0.3), "上海创智汇 × 同浦汇｜活动专项交付", size=14, color=BLUE, bold=True)
    text(s, ML, Inches(1.95), Inches(12), Inches(0.7), "创智汇 · 30场活动怎么做", size=34, color=INK, bold=True)
    text(s, ML, Inches(2.75), Inches(12), Inches(0.35), "2026.08—2027.07　｜　同步衔接 WAIC 与 ChinaJoy", size=15, color=GOLD, bold=True)
    text(s, ML, Inches(3.2), Inches(12), Inches(0.35), "每场≤30人　·　负责人占比≤30%　·　我们带客 / 园区销售", size=13, color=MUT)
    card(s, ML, Inches(3.8), Inches(6.0), Inches(2.85), "主体与支持单位", [
        "运营主体范围：上海市云计算创新基地（国家级孵化器）",
        "学术支持：复旦大学住房政策研究中心",
        "活动及运营载体支持：杨浦区科技企业联合会、科技企业服务中心",
        "执行：同浦汇（策划执行与带客）",
    ], BLUE)
    card(s, Inches(7.2), Inches(3.8), Inches(5.55), Inches(2.85), "本册边界", [
        "仅活动专项交付，不做租金租赁对赌",
        "因租金3.3元/㎡/天高于周边均价，不做租赁必要性要求",
        "出海活动另议另计价，不在本册费用内",
        "领事到访 / 挂牌活动另计价（三部分费用之外）",
        "不承诺一定带来外资企业",
    ], GREEN)
    footer(s, 1, TOTAL)

    # 2 主体单位
    s = slide(prs); header(s, "01 · 主体与支持", "谁在做这件事")
    rows = [
        ["角色", "单位", "说明"],
        ["主体运营范围", "上海市云计算创新基地", "国家级孵化器，纳入创智汇活动与运营主体范围"],
        ["学术支持", "复旦大学住房政策研究中心", "学术背书、议题与人才/空间相关研究支持"],
        ["活动及运营载体支持", "杨浦区科技企业联合会", "活动生态、会员协同、传播支持"],
        ["活动及运营载体支持", "杨浦区科技企业服务中心", "政策服务、申报辅导、入驻手续协同"],
        ["活动执行与带客", "同浦汇", "策划、邀约、执行、建群、带客回访"],
        ["销售成交", "园区", "房源销售、价格确认、合同交割"],
    ]
    table(s, Inches(1.35), rows, [0.22, 0.32, 0.46], sizes=[12, 14, 13], rh=0.55)
    footer(s, 2, TOTAL)

    # 3 口径
    s = slide(prs); header(s, "01 · 关键口径", "人数 · 负责人 · 客群 · 外企原则")
    card(s, ML, Inches(1.35), Inches(4.0), Inches(5.4), "人数与负责人", [
        "每场到场人次控制在 30 人以下",
        "全年触达目标约 600 人次+，以可执行、可核验为准",
        "负责人占比≤30%",
        "邀约以园区产业主题客群为主（弱化关联度表述）",
    ], BLUE)
    card(s, Inches(4.75), Inches(1.35), Inches(4.0), Inches(5.4), "外企与领事", [
        "原则：不承诺一定会带外资企业来",
        "出海活动：另议、价格另议，不在本册费用内",
        "领事到访接待：另计价",
        "领事馆挂牌及挂牌活动：另计价",
        "挂牌与领事活动性质不同，分开计价",
    ], GOLD)
    card(s, Inches(8.95), Inches(1.35), Inches(3.8), Inches(5.4), "分工一句话", [
        "我们负责带客",
        "园区负责销售部分",
        "活动结束建群",
        "园区可在群内促成成交",
    ], GREEN)
    footer(s, 3, TOTAL)

    # 4-5 总表
    for subset, title, pg in ((EVENTS[:15], "30场概况总表（上）", 4), (EVENTS[15:], "30场概况总表（下）", 5)):
        s = slide(prs); header(s, "02 · 活动概况", title, "主题 · 时间 · 形式 · 人数≤30 · 负责人≤30% · 价值落点")
        rows = [["编号", "主题/活动", "时间", "形式", "人数", "负责人", "价值落点"]]
        for e in subset:
            rows.append([e["code"], e["name"][:16], e["month"].replace("2026年", "").replace("2027年", "次"),
                         e["fmt"][:8], e["people"], e["dm"], e["lease"][:14]])
        table(s, Inches(1.25), rows, [0.08, 0.3, 0.14, 0.12, 0.1, 0.1, 0.16], sizes=[9, 10, 9, 9, 9, 9, 9], rh=0.32)
        footer(s, pg, TOTAL)

    # 6 价值
    s = slide(prs); header(s, "03 · 活动价值", "活动带来什么")
    card(s, ML, Inches(1.35), Inches(5.9), Inches(5.4), "园区侧", [
        "每月有主题不断档，空间有热度",
        "8月项目推广日打通中介/银行/投研/政府",
        "内容场次对齐 ChinaJoy 数字娱乐主题",
        "全年节奏同步衔接 WAIC 与 ChinaJoy",
        "云创基地国家级孵化器背书可感知",
    ], BLUE)
    card(s, Inches(7.2), Inches(1.35), Inches(5.55), Inches(5.4), "转化侧", [
        "每场≤30人，精而准",
        "固定看场环节，便于带客",
        "活动结束建群，持续触达",
        "同浦汇带客，园区做销售促成",
        "不做租金租赁对赌与必要性要求",
    ], GREEN)
    footer(s, 6, TOTAL)

    # 7 交付
    s = slide(prs); header(s, "03 · 具体交付", "我们交给园区的东西")
    rows = [
        ["周期", "交付物", "说明"],
        ["签约后2周", "年度活动方案细化版 + 排期日历", "活动专项"],
        ["每场结束后7天", "执行卡、议程、嘉宾确认、报名页、签到表、意向表、动线图、照片/通稿", "同浦汇归档"],
        ["每场结束后7天", "回访摘要（知会版）+ 建群清单", "便于园区销售跟进与群内促成"],
        ["每月", "活动数据月报 + 下月排期确认单", "联席会前"],
        ["每季/年末", "季度复盘、全年台账、WAIC/ChinaJoy衔接建议", "资产沉淀"],
    ]
    table(s, Inches(1.35), rows, [0.18, 0.52, 0.3], sizes=[12, 12, 12], rh=0.58)
    text(s, ML, Inches(5.9), CW, Inches(0.7),
         "回访与带客由同浦汇负责；园区聚焦销售促成。交园区的回访材料为摘要知会，方便销售跟进，不额外增加园区事务负担。",
         size=12, color=MUT)
    footer(s, 7, TOTAL)

    # 8 分工
    s = slide(prs); header(s, "03 · 执行分工", "我们带客，园区负责销售")
    rows = [
        ["事项", "同浦汇", "园区", "载体支持/学术"],
        ["策划与议程", "全案策划、执行卡", "方案审定、档期确认", "联合会/服中心协同"],
        ["邀约与带客", "定向邀约、带客到场、建群", "共享在谈名单", "会员/学者推荐"],
        ["场地保障", "提前提报需求、物料", "报备、导视、物业安保、样板间", "—"],
        ["现场执行", "主持控场、看场带队", "现场支持、开放看场", "政策宣讲位（如需）"],
        ["销售促成", "提供意向名单与群", "负责销售部分、报价成交、合同", "入驻手续（服中心）"],
        ["回访跟踪", "主责带客回访与台账", "基于名单/群做销售跟进", "—"],
        ["复盘", "单场复盘+月报", "确认下月活动安排", "—"],
    ]
    table(s, Inches(1.3), rows, [0.14, 0.3, 0.3, 0.26], sizes=[11, 11, 11, 11], rh=0.5)
    footer(s, 8, TOTAL)

    # 9 协同
    s = slide(prs); header(s, "03 · 协同着眼点", "做完活动、建完群 → 园区销售促成")
    card(s, ML, Inches(1.35), Inches(5.9), Inches(5.4), "标准动作", [
        "D日：同浦汇签到分层并建群",
        "T+1：同浦汇发看场预约与资料入群",
        "T+1~7：同浦汇持续带客回访",
        "园区：在群内发布房源/政策/条件，负责销售促成",
        "双方月度核对意向与成交进展",
    ], GREEN)
    card(s, Inches(7.2), Inches(1.35), Inches(5.55), Inches(5.4), "责任边界", [
        "我们负责带客",
        "园区负责销售部分",
        "回访摘要交园区便于销售跟进",
        "不承诺必带外资企业",
        "出海/领事相关另议另计价",
    ], BLUE)
    footer(s, 9, TOTAL)

    # 10 收费付款
    s = slide(prs); header(s, "04 · 收费与付款", "活动费用边界 · 付款节点（唯一方案）")
    rows = [
        ["项目", "口径", "说明"],
        ["活动策划执行服务费", "年度打包（金额线下确认）", "含常规30场策划执行、物料、主持、建群与带客回访"],
        ["租金参考", "3.3 元/㎡/天", "高于周边均价，故不做租赁对赌与租赁必要性要求"],
        ["建议免租期", "建议 1–3 个月（面议）", "降低决策门槛"],
        ["物业参考", "13.8 元/㎡/月", "园区统一标准"],
        ["待租房源", "8 间优先看场", "带客当场可看"],
        ["出海活动", "另议 · 另计价", "不在本册活动服务费内"],
        ["领事到访接待", "另计价", "在现有三部分费用之外"],
        ["领事馆挂牌及挂牌活动", "另计价", "与领事到访分属不同性质，分开计价"],
    ]
    table(s, Inches(1.2), rows, [0.24, 0.28, 0.48], sizes=[11, 12, 11], rh=0.38)
    text(s, ML, Inches(5.35), CW, Inches(0.28), "付款节点（唯一建议方案）", size=13, color=BLUE, bold=True)
    rows2 = [
        ["节点", "比例", "对应"],
        ["签约后7日内", "50% 启动款", "启动策划、年度日历锁定"],
        ["Q3 / Q4 / 次年Q1 / 次年Q2 各季", "各10%（合计40%）", "与季度场次完成度、月报验收挂钩"],
        ["年终收官后", "10% 尾款", "年报、台账、影像资产移交"],
    ]
    table(s, Inches(5.65), rows2, [0.4, 0.28, 0.32], sizes=[11, 11, 11], rh=0.34)
    footer(s, 10, TOTAL)

    # 11 执行节奏
    s = slide(prs); header(s, "05 · 单场执行节奏", "同一套打法")
    rows = [
        ["节点", "动作", "责任", "交付"],
        ["T-14", "锁定嘉宾议程物料", "同浦汇", "议程/嘉宾确认/报名"],
        ["T-7", "邀约确认+场地报备", "同浦汇；园区场地确认", "名单/场地单"],
        ["T-1", "彩排+看场动线", "同浦汇；园区样板间", "彩排记录"],
        ["D日", "签到→主区→看场→建群", "同浦汇执行；园区支持", "签到/意向/群"],
        ["T+1~7", "带客回访", "同浦汇主责", "回访摘要知会园区"],
        ["销售促成", "房源条件与成交", "园区负责销售部分", "合同/交割"],
    ]
    table(s, Inches(1.35), rows, [0.12, 0.3, 0.32, 0.26], sizes=[11, 12, 11, 11], rh=0.55)
    footer(s, 11, TOTAL)

    # 12-13 cases
    for i, c in enumerate(CASES):
        s = slide(prs); header(s, f"06 · 往期案例（{i+1}）", c["name"], c["meta"])
        for j, fn in enumerate(c["photos"][:3]):
            add_pic(s, os.path.join(PHOTO, fn), Emu(ML + j * Inches(4.1)), Inches(1.35), Inches(3.95), Inches(2.3))
        rect(s, ML, Inches(3.85), Inches(7.5), Inches(2.95), fill=BG, line=LINE)
        text(s, Emu(ML + Inches(0.2)), Inches(3.95), Inches(7.1), Inches(0.3), "当天议程", size=13, color=BLUE, bold=True)
        bullets(s, Emu(ML + Inches(0.2)), Inches(4.35), Inches(7.1), Inches(2.2), c["agenda"], size=13)
        rect(s, Inches(8.3), Inches(3.85), Inches(4.5), Inches(2.95), fill=ACC, line=LINE)
        text(s, Inches(8.45), Inches(3.95), Inches(4.2), Inches(0.3), "可见机构/企业", size=13, color=BLUE, bold=True)
        text(s, Inches(8.45), Inches(4.35), Inches(4.2), Inches(2.2),
             f"{c['orgs']}\n\n企业侧：{c['companies']}\n\n映射：{c['map']}", size=12, color=INK)
        footer(s, 12 + i, TOTAL)

    # 14-16 逐场
    chunks = [(EVENTS[:10], 14), (EVENTS[10:20], 15), (EVENTS[20:], 16)]
    for subset, pg in chunks:
        s = slide(prs); header(s, "07 · 逐场执行要素", f"{subset[0]['code']}–{subset[-1]['code']}｜行程/嘉宾/交付")
        rows = [["编号", "活动", "行程", "嘉宾", "交付"]]
        for e in subset:
            rows.append([e["code"], e["name"][:14], e["agenda"][:24], e["guests"][:14], "标准包+建群" if "另" not in e["deliver"] else "另案"])
        table(s, Inches(1.3), rows, [0.07, 0.22, 0.3, 0.26, 0.15], sizes=[9, 10, 10, 10, 10], rh=0.42)
        footer(s, pg, TOTAL)

    # 17 ChinaJoy + WAIC
    s = slide(prs); header(s, "08 · 双会衔接", "活动同步衔接 WAIC 与 ChinaJoy")
    card(s, ML, Inches(1.35), Inches(5.9), Inches(5.4), "WAIC 线", [
        "8月开放交流日承接大会议题余热",
        "Agent/多模态/算力/具身训练营常态化",
        "年底 Demo Day 收口",
        "次年创作者首发预热 WAIC 2027",
        "云创基地算力与孵化资源贯穿",
    ], BLUE)
    card(s, Inches(7.2), Inches(1.35), Inches(5.55), Inches(5.4), "ChinaJoy 线", [
        "数字娱乐与IP授权对接",
        "游戏周边与潮玩供应链对接",
        "创作者内容首发",
        "数字娱乐市集体验 / 沉浸式联展",
        "IP联名与衍生品撮合",
    ], GOLD)
    footer(s, 17, TOTAL)

    # 18 另计价说明
    s = slide(prs); header(s, "08 · 另议另计价事项", "不在本册活动服务费内")
    rows = [
        ["事项", "计价方式", "性质说明"],
        ["出海专题活动", "另议 · 另计价", "不纳入本册年度活动包"],
        ["领事到访接待", "另计价", "外事接待类"],
        ["领事馆挂牌及挂牌活动", "另计价", "挂牌类；与领事到访分属不同性质"],
        ["外资企业到场", "不承诺", "原则：不承诺一定会带外资企业来"],
    ]
    table(s, Inches(1.4), rows, [0.28, 0.28, 0.44], sizes=[13, 13, 13], rh=0.7)
    footer(s, 18, TOTAL)

    # 19 结论
    s = slide(prs); header(s, "09 · 一页纸结论", "活动专项交付要点")
    card(s, ML, Inches(1.35), Inches(4.0), Inches(5.4), "我们交付", [
        "30场排期（每场≤30人）",
        "8月项目推广日（中介/五大行/投研/政府）",
        "ChinaJoy主题内容场次",
        "同步衔接 WAIC 与 ChinaJoy",
        "带客、建群、回访摘要",
        "往期案例（企业+议程+现场照）",
    ], BLUE)
    card(s, Inches(4.75), Inches(1.35), Inches(4.0), Inches(5.4), "请园区协同", [
        "场地物业与样板间",
        "共享在谈名单、确认档期",
        "负责销售部分与成交交割",
        "群内发布条件促成成交",
        "云创基地/联合会/服中心协同到位",
    ], GREEN)
    card(s, Inches(8.95), Inches(1.35), Inches(3.8), Inches(5.4), "商务口径", [
        "付款：50%+4×10%+10%",
        "租金不做对赌/必要性要求",
        "建议免租期1–3个月",
        "出海/领事另计价",
        "不承诺必带外企",
    ], GOLD)
    footer(s, 19, TOTAL)

    # 20 end
    s = slide(prs)
    text(s, ML, Inches(2.3), Inches(12), Inches(0.5), "谢谢审阅", size=36, color=INK, bold=True, align=PP_ALIGN.CENTER)
    text(s, ML, Inches(3.1), Inches(12), Inches(0.4), "创智汇 · 30场活动专项交付", size=18, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    text(s, ML, Inches(3.7), Inches(12), Inches(0.35), "我们带客 · 园区销售 · 衔接 WAIC & ChinaJoy", size=14, color=MUT, align=PP_ALIGN.CENTER)
    text(s, ML, Inches(4.4), Inches(12), Inches(0.9),
         "主体运营：上海市云计算创新基地（国家级孵化器）\n学术支持：复旦大学住房政策研究中心\n载体支持：杨浦区科技企业联合会 · 科技企业服务中心",
         size=13, color=SOFT, align=PP_ALIGN.CENTER)
    footer(s, 20, TOTAL)

    os.makedirs(ART, exist_ok=True)
    outs = [
        "创智汇30场活动专项交付-修订版.pptx",
        "chuangzhihui-30-events-activity-only-rev.pptx",
        "创智汇30场活动具体方案.pptx",
        "chuangzhihui-30-events-plan.pptx",
    ]
    tmp = os.path.join(HERE, "_tmp_act_v2.pptx")
    prs.save(tmp)
    for fn in outs:
        p = os.path.join(HERE, fn)
        shutil.copy2(tmp, p)
        shutil.copy2(p, os.path.join(ART, fn))
    os.remove(tmp)
    print("PPT slides:", len(prs.slides))

    # Excel
    wb = Workbook()
    thin = Border(left=Side(style="thin", color="D0D0D0"), right=Side(style="thin", color="D0D0D0"),
                  top=Side(style="thin", color="D0D0D0"), bottom=Side(style="thin", color="D0D0D0"))
    head = PatternFill("solid", fgColor="1E4D8C")
    alt = PatternFill("solid", fgColor="EEF2F8")

    def hdr(ws, headers):
        for j, h in enumerate(headers, 1):
            c = ws.cell(1, j, h); c.font = Font(name="微软雅黑", bold=True, color="FFFFFF", size=10)
            c.fill = head; c.alignment = Alignment(wrap_text=True, vertical="center"); c.border = thin

    def write(ws, data, widths, h=44):
        for i, row in enumerate(data, 2):
            for j, v in enumerate(row, 1):
                c = ws.cell(i, j, v); c.font = Font(name="微软雅黑", size=9)
                c.alignment = Alignment(wrap_text=True, vertical="center"); c.border = thin
                if i % 2 == 0: c.fill = alt
            ws.row_dimensions[i].height = h
        ws.row_dimensions[1].height = 28
        for i, w in enumerate(widths, 1):
            ws.column_dimensions[get_column_letter(i)].width = w

    ws = wb.active; ws.title = "01-活动概况"
    hdr(ws, ["编号", "活动名称", "时间", "形式", "人数", "负责人占比", "场地", "节奏", "价值落点", "行程", "嘉宾", "交付"])
    write(ws, [[e["code"], e["name"], e["month"], e["fmt"], e["people"], e["dm"], e["floor"], e["rhythm"],
                e["lease"], e["agenda"], e["guests"], e["deliver"]] for e in EVENTS],
          [8, 30, 16, 12, 10, 10, 14, 28, 22, 36, 28, 28], 50)

    ws2 = wb.create_sheet("02-主体分工收费付款")
    hdr(ws2, ["模块", "内容"])
    write(ws2, [
        ["主体运营范围", PARTIES["ops"]],
        ["学术支持", PARTIES["academic"]],
        ["活动及运营载体支持", PARTIES["support"]],
        ["执行带客", PARTIES["exec"]],
        ["销售", "园区负责销售部分与合同交割"],
        ["人数", "每场≤30人；全年触达约600人次+，可执行可核验"],
        ["负责人占比", "≤30%"],
        ["客群", "以园区产业主题邀约为主（弱化关联度/核验表表述）"],
        ["外企原则", "不承诺一定会带外资企业来"],
        ["出海", "另议另计价，不在本册费用内"],
        ["领事到访", "另计价（三部分费用之外）"],
        ["领事挂牌/挂牌活动", "另计价；与到访分属不同性质"],
        ["租金", "3.3元/㎡/天高于周边均价，不做租赁对赌与必要性要求"],
        ["免租期建议", "1–3个月（面议）"],
        ["付款", "签约后7日50%；Q3/Q4/次年Q1/次年Q2各10%；年终10%"],
        ["双会衔接", "活动同步衔接 WAIC 与 ChinaJoy"],
    ], [18, 72], 34)

    ws3 = wb.create_sheet("03-往期案例")
    hdr(ws3, ["案例", "时间地点", "机构", "企业", "议程", "照片", "映射"])
    write(ws3, [[c["name"], c["meta"], c["orgs"], c["companies"], " | ".join(c["agenda"]), "；".join(c["photos"]), c["map"]] for c in CASES],
          [32, 28, 28, 36, 40, 28, 24], 56)

    for fn in ["创智汇30场活动专项交付-修订版.xlsx", "chuangzhihui-30-events-activity-only-rev.xlsx",
               "创智汇30场活动方案排期表.xlsx", "chuangzhihui-30-events-schedule.xlsx"]:
        p = os.path.join(HERE, fn); wb.save(p); shutil.copy2(p, os.path.join(ART, fn))
    print("Excel OK")


if __name__ == "__main__":
    build()
