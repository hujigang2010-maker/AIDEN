# -*- coding: utf-8 -*-
"""创智汇 · 30场活动方案（WAIC 2026议程对齐升级版）
结合 WAIC 2026「智能伙伴 共创未来」+ WAICA学术议程 + 北外滩双夜逻辑，
升级创智汇全年30场活动（总—分—总）。
用法: python3 build_events_30_waic.py
"""
import os, shutil
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
from openpyxl.utils import get_column_letter

HERE = os.path.dirname(os.path.abspath(__file__))
ART = "/opt/cursor/artifacts"
FONT, FONT_EN = "Microsoft YaHei", "Arial"
BG_A, BG_B = "1A1140", "3A206E"
INK = RGBColor(0xF5, 0xF1, 0xFF)
MUT = RGBColor(0xC6, 0xBA, 0xEA)
SOFT = RGBColor(0x94, 0x86, 0xC4)
ACC = RGBColor(0x8B, 0x7B, 0xFF)
ACC2 = RGBColor(0xC7, 0x7D, 0xFF)
GOLD = RGBColor(0xE9, 0xC2, 0x7C)
GREEN = RGBColor(0x56, 0xD6, 0xB6)
ROSE = RGBColor(0xF5, 0x8B, 0xB8)
LINE = RGBColor(0x4A, 0x3C, 0x7C)
SW, SH = Inches(13.333), Inches(7.5)
ML, CW = Inches(0.75), Inches(11.83)
prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
BLANK = prs.slide_layouts[6]; _FOOT = []; IDX = 0

# WAIC 2026 议程锚点（公开议程 + 学术日程摘要）
WAIC_META = {
    "theme": "智能伙伴 共创未来",
    "when": "2026.7.17–7.20",
    "where": "浦东世博 · 张江 · 徐汇西岸「三地四馆」",
    "tracks": [
        ("开幕主旨", "智能伙伴/全球AI治理/图灵·诺奖嘉宾"),
        ("Agent智能体", "Agent工具评测、工作流、多代理辩论"),
        ("AI for Science", "科学多模态智能体、数学建模与科学计算"),
        ("算力Infra", "GPU集群调度、端侧推理、投机解码"),
        ("具身智能", "空间交互、机器人失效分析与矫正"),
        ("青年菁英", "WAICA青年菁英会、Builders建造者"),
        ("AGI思辨", "模型下一步、安全对齐与治理"),
        ("内容创作", "官方创作者之夜错位协同·园区IP落地"),
    ],
}

# cat: A智能体训练 B政策 C产业对接 D内容IP E WAIC大场 F出海/AGI
EVENTS = [
    # A ×8 —— 对齐 Agent / 多模态 / 算力 / 具身
    dict(id=1, cat="A", waic="智能伙伴", name="智能伙伴·创智汇开年主题沙龙", month="1月下旬", floor="5F沙龙",
         dur="2.5小时", guests="联合会领导、AI企业代表、服务中心",
         agenda=["18:30 签到参观", "19:00 WAIC主题年度解读", "19:30 园区AI+内容定位发布", "20:10 圆桌：智能伙伴如何落杨浦", "20:40 看场+意向"],
         park="把WAIC年度主题转化为园区主叙事", lease="开年集中获客与品牌定调", format="主题沙龙", staff="我方3+园区2+联合会1"),
    dict(id=2, cat="A", waic="Agent", name="Agent智能体搭建一日营", month="2月中旬", floor="3F OPC",
         dur="1天", guests="Agent工程师、开源Maintainer",
         agenda=["09:30 签到", "10:00 Agent工作流拆解", "13:30 低代码实操（对齐MCP工具调用）", "16:00 Demo路演", "16:40 看场逼单"],
         park="对齐WAICA Agent/工具评测赛道", lease="Agent初创工位/单元转化", format="实训营", staff="我方3+园区2"),
    dict(id=3, cat="A", waic="多模态", name="多模态智能体工作坊", month="3月上旬", floor="3F培训+直播间",
         dur="1天", guests="多模态研究者、应用厂商",
         agenda=["多模态案例", "视觉-语言实操", "行业场景共创", "作品评审", "政策礼包+看场"],
         park="承接WAICA「科学多模态智能体」议题落地", lease="多模态应用团队入驻", format="工作坊", staff="我方4+园区2"),
    dict(id=4, cat="A", waic="算力Infra", name="火山引擎×算力Infra实务营", month="3月下旬", floor="3F+云创基地",
         dur="1天", guests="火山引擎讲师、云计算创新基地",
         agenda=["算力政策包", "代金券/折扣实操", "集群调度案例（呼应WAIC Infra）", "一对一诊断", "入驻捆绑云资源"],
         park="国级孵化器+专有算力包兑现", lease="高算力消耗型企业定向招商", format="厂商联训", staff="我方3+园区2+厂商2"),
    dict(id=5, cat="A", waic="具身智能", name="具身智能空间交互体验日", month="4月上旬", floor="5F展区角+3F",
         dur="半天", guests="机器人/具身团队、高校实验室",
         agenda=["具身演示", "空间交互讲解", "场景清单对接", "5F应用场景踩点", "落位洽谈"],
         park="对齐WAICA具身智能空间交互工作坊", lease="具身/机器人团队看场", format="体验日", staff="我方3+园区2"),
    dict(id=6, cat="A", waic="AI内容", name="AIGC微短剧制片特训", month="5月上旬", floor="3F直播间+5F",
         dur="1天", guests="微短剧导演、平台方",
         agenda=["YOUNG立方政策", "AIGC制片流程", "脚本工作坊", "优秀作品路演", "5F场景看场"],
         park="内容赛道+官方创作者夜场错位协同", lease="厂牌/工作室入驻", format="特训营", staff="我方4+园区2"),
    dict(id=7, cat="A", waic="Builders", name="OPC超级个体黑客松（春）", month="5月中旬", floor="3F整层",
         dur="1.5天", guests="评委×3、投资人×2、云厂商",
         agenda=["开题组队", "封闭开发", "路演决赛", "入驻礼包", "一对一看场"],
         park="青年菁英/Builders氛围", lease="获奖团队优先谈单元", format="黑客松", staff="我方5+园区3"),
    dict(id=8, cat="A", waic="Agent", name="AI营销Agent实战营", month="10月中旬", floor="3F培训",
         dur="半天", guests="投放操盘手、Agent产品经理",
         agenda=["投放策略", "Agent自动化", "素材AIGC", "复盘", "服务包转化"],
         park="企业服务粘性", lease="营销型公司入驻", format="实战营", staff="我方3+园区1"),
    # B ×5 政策
    dict(id=9, cat="B", waic="智能伙伴", name="YOUNG立方×智能伙伴政策沙龙", month="1月中旬", floor="5F沙龙",
         dur="2.5小时", guests="区政策宣讲、服务中心",
         agenda=["政策要点", "WAIC主题与园区礼包", "适用画像", "诊断预约", "看场"],
         park="政策前置+主题统一", lease="内容/AI企业导入", format="沙龙", staff="我方2+园区2+服务中心2"),
    dict(id=10, cat="B", waic="治理", name="AI治理与可信智能体沙龙", month="2月下旬", floor="3F",
         dur="2小时", guests="治理学者、合规顾问、企业法务",
         agenda=["全球治理框架速递", "企业合规清单", "案例", "问答", "入驻服务对接"],
         park="呼应WAIC全球AI治理高级别会议", lease="合规意识型企业信任入驻", format="沙龙", staff="我方2+服务中心2"),
    dict(id=11, cat="B", waic="产业", name="高企认定冲刺（AI企业专场）", month="3月中旬", floor="3F",
         dur="半天", guests="高企辅导、财税顾问",
         agenda=["条件拆解", "AI企业材料要点", "初筛", "辅导套餐", "入驻激励"],
         park="在园资质提升", lease="外区待认定AI企业带政策入驻", format="辅导会", staff="我方2+服务中心3"),
    dict(id=12, cat="B", waic="产业", name="专精特新·AI应用培育路演", month="6月上旬", floor="5F",
         dur="半天", guests="评审顾问、银行、基金",
         agenda=["梯度政策", "企业路演", "点评", "金融对接", "看场"],
         park="优质企业筛选", lease="成长型AI企业补位", format="路演", staff="我方3+园区2"),
    dict(id=13, cat="B", waic="算力", name="创新券·算力券·模型券实务沙龙", month="11月上旬", floor="3F",
         dur="2小时", guests="创新券平台、云厂商",
         agenda=["三券规则", "核销实操", "案例", "开户引导", "入驻转化"],
         park="服务机构KPI与核销流水", lease="用券企业向园区聚集", format="实务沙龙", staff="我方2+服务中心2"),
    # C ×5 对接
    dict(id=14, cat="C", waic="AI for Science", name="高校成果转化·AI for Science日", month="4月中旬", floor="3F+沙龙",
         dur="半天", guests="复旦/同济技术转移、教授团队",
         agenda=["成果路演", "科学智能体场景", "园区承接方案", "看孵化单元", "转化落位"],
         park="对齐WAICA科学多模态/AI for Science", lease="成果公司/实验室落户", format="对接日", staff="我方3+园区2+高校2"),
    dict(id=15, cat="C", waic="产业", name="汕头玩具×AIGC供应链对接会", month="4月下旬", floor="5F玩具区",
         dur="半天", guests="汕头商协会、品牌采购",
         agenda=["集群介绍", "AIGC设计提效", "供需对接", "展位参观", "报价"],
         park="5F集群填充+AI赋能传统产业", lease="玩具企业展位去化", format="对接会", staff="我方3+园区2"),
    dict(id=16, cat="C", waic="产业", name="扬州毛绒×数字人联名对接", month="6月中旬", floor="5F毛绒区",
         dur="半天", guests="扬州集群、数字人厂商",
         agenda=["联名模式", "租金/扣点测算", "看场", "MOU", "跟进"],
         park="毛绒集群落地", lease="700㎡量级客户", format="对接会", staff="我方3+园区2"),
    dict(id=17, cat="C", waic="产业", name="东莞潮玩品牌入沪推介", month="8月上旬", floor="5F潮玩区",
         dur="半天", guests="潮玩品牌、渠道商",
         agenda=["上海渠道", "5F落位", "联名活动", "看场", "定金意向"],
         park="潮玩内容填充", lease="品牌展位/快闪", format="推介会", staff="我方3+园区2"),
    dict(id=18, cat="C", waic="内容", name="IP授权×AI衍生交易撮合会", month="9月中旬", floor="5F展示中心",
         dur="1天", guests="IP方、被授权商、律师",
         agenda=["IP路演", "AI衍生专题桌", "一对一撮合", "成交看板", "落户激励"],
         park="IP交易+AI衍生", lease="IP/衍生品团队入驻", format="撮合会", staff="我方4+园区2"),
    # D ×4 内容IP
    dict(id=19, cat="D", waic="内容", name="AI创作者内容首发①", month="5月下旬", floor="5F主展",
         dur="3小时", guests="创作者、媒体、渠道",
         agenda=["揭幕", "发布", "签售快闪", "媒体专访", "招商通道"],
         park="与官方创作者之夜错位：园区做「可落位」首发", lease="关联品牌问询", format="发布会", staff="我方4+园区3"),
    dict(id=20, cat="D", waic="内容", name="潮玩×AIGC主题市集", month="6月下旬", floor="5F+公区",
         dur="2天", guests="摊主×40、达人",
         agenda=["布摊", "开市", "AIGC互动打卡", "夜间场", "优质摊主转正谈"],
         park="人气与经营收入", lease="摊主升级固定展位/办公", format="市集", staff="我方5+园区4"),
    dict(id=21, cat="D", waic="内容", name="沉浸式AI+IP联展", month="8月中旬", floor="5F展示中心",
         dur="7–10天", guests="联合IP×5、文旅渠道",
         agenda=["布展", "开幕", "预约观展", "教育场", "闭幕招商酒会"],
         park="长期人流", lease="闭幕酒会集中转化", format="联展", staff="我方4+园区3+外包"),
    dict(id=22, cat="D", waic="内容", name="AI创作者内容首发②", month="10月下旬", floor="5F",
         dur="3小时", guests="IP方、渠道、媒体",
         agenda=["发布", "双十一预售", "渠道对接", "看场", "年框意向"],
         park="下半年内容高峰", lease="旺季补位", format="发布会", staff="我方4+园区3"),
    # E ×4 WAIC大场（核心升级）
    dict(id=23, cat="E", waic="WAIC主会", name="WAIC UP! 创智汇「智能伙伴」开放日", month="7月17日", floor="3F+5F整区",
         dur="1天（10:00–17:00）", guests="参展回流嘉宾、媒体、企业、高校",
         agenda=["10:00 开幕致辞（联合会+同浦汇）", "10:30 WAIC主题速递与园区承接", "11:30 空间开放参观", "14:00 Agent/内容双展台", "15:30 政策礼包发布", "16:30 意向洽谈"],
         park="WAIC日间散场承接、杨浦可见度", lease="大会流量导入园区看场", format="开放日", staff="我方6+园区5+服务中心2"),
    dict(id=24, cat="E", waic="Builders夜", name="WAIC UP! AI Builders Night·建造者之夜", month="7月17日晚", floor="5F主厅+露台/沙龙",
         dur="18:30–21:30", guests="创业者、开发者、投资人、开源Maintainer",
         agenda=["18:30 签到·名牌「我正在造的AI」", "18:40 致辞", "18:50 闪电演讲×4", "19:40 1分钟开放麦", "20:10 AMA", "20:40 投资人蹲项目区+看场"],
         park="对齐北外滩双夜Day1逻辑，落位创智汇", lease="建造者当场看3F单元", format="夜场社交", staff="我方5+园区3+联合会2"),
    dict(id=25, cat="E", waic="AGI夜", name="WAIC UP! 通往AGI之夜", month="7月18日晚", floor="5F圆桌+学术酒吧分区",
         dur="19:00–21:30", guests="复旦学者、大模型负责人、Agent创始人、治理学者",
         agenda=["19:00 复旦致辞定调", "19:10 无PPT圆桌：通往AGI走到哪了", "20:10 三桌围炉（模型/智能体/科学与治理）", "20:40 辩题投票", "21:00 深度社交+精品看场"],
         park="对齐WAICA思辨+复旦学术支持", lease="高端涉外/研究型企业入驻", format="学术夜场", staff="我方4+园区2+复旦2"),
    dict(id=26, cat="E", waic="青年菁英", name="创智汇秋季潮玩×AI嘉年华", month="10月上旬", floor="5F+公区",
         dur="2天", guests="潮玩品牌、达人、赞助商、AI创作者",
         agenda=["开幕秀", "市集+展洽", "AI互动赛事", "夜间场", "签约仪式"],
         park="秋季品牌大事件、WAIC精神延续", lease="集中签约展位/办公", format="嘉年华", staff="我方8+园区6+外包"),
    # F ×4 出海/AGI/年终
    dict(id=27, cat="F", waic="出海", name="AI出海推介与国际买家团", month="9月下旬", floor="5F+培训",
         dur="1天", guests="海外买家、跨境平台、出海顾问",
         agenda=["出海路径", "买家需求", "一对一洽谈", "订单意向", "办公/展位捆绑"],
         park="出海服务闭环", lease="出海需求企业入驻", format="买家团", staff="我方5+园区2"),
    dict(id=28, cat="F", waic="AGI", name="通往AGI季度圆桌①", month="4月下旬", floor="5F沙龙",
         dur="2.5小时", guests="学者、模型企业、投资人",
         agenda=["议题导入", "圆桌", "围炉", "问答", "看场"],
         park="AGI议题全年运营（WAIC会后延续）", lease="研究型/模型团队", format="圆桌", staff="我方3+园区1+学术1"),
    dict(id=29, cat="F", waic="领事", name="国别会客厅·AI合作领事专题", month="11月中旬", floor="会客厅/5F",
         dur="2.5小时", guests="领事官员、涉外AI企业、翻译",
         agenda=["外事接待", "国别AI机会", "B2B闪见", "国际服务介绍", "高端看场"],
         park="国际背书", lease="涉外/出海企业", format="领事沙龙", staff="我方4+园区2+外事"),
    dict(id=30, cat="F", waic="智能伙伴", name="创智汇AI年度Demo Day·智能伙伴收官", month="12月上旬", floor="5F主展",
         dur="1天", guests="投资人、链主、媒体、区级嘉宾",
         agenda=["开幕", "十强路演", "颁奖", "明年WAIC预热", "红酒洽谈+集中签约"],
         park="年度成果+来年WAIC预热", lease="集中签约与媒体背书", format="路演日", staff="我方6+园区4"),
]

CAT_META = {
    "A": ("智能体/多模态/算力训练", "8场", "对齐Agent·科学智能体·Infra·具身", GOLD),
    "B": ("政策与治理沙龙", "5场", "对齐智能伙伴·治理·三券", ACC),
    "C": ("产业与科学对接", "5场", "对齐AI for Science·集群", ACC2),
    "D": ("AI内容与IP", "4场", "与官方创作者夜场错位协同", GREEN),
    "E": ("WAIC联动大场", "4场", "7月开放日+双夜+秋季嘉年华", ROSE),
    "F": ("出海/AGI/收官", "4场", "会后延续与国际链接", GOLD),
}


def _grad(shape, stops, ang=90):
    spPr = shape._element.spPr
    for t in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill"):
        e = spPr.find(qn(t))
        if e is not None: spPr.remove(e)
    g = spPr.makeelement(qn("a:gradFill"), {}); lst = g.makeelement(qn("a:gsLst"), {})
    for pos, col, al in stops:
        gs = g.makeelement(qn("a:gs"), {"pos": str(int(pos * 1000))})
        c = g.makeelement(qn("a:srgbClr"), {"val": col})
        if al is not None: c.append(g.makeelement(qn("a:alpha"), {"val": str(int(al * 1000))}))
        gs.append(c); lst.append(gs)
    g.append(lst); g.append(g.makeelement(qn("a:lin"), {"ang": str(int(ang * 60000)), "scaled": "1"}))
    ln = spPr.find(qn("a:ln")); (ln.addprevious(g) if ln is not None else spPr.append(g))


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.shadow.inherit = False; bg.line.fill.background()
    _grad(bg, [(0, BG_A, None), (55, BG_B, None), (100, BG_A, None)], 120)
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s


def rect(s, x, y, w, h, fill=None, line=None, lw=1.0, radius=False, grad=None, gang=90):
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    b.shadow.inherit = False
    if grad is not None: _grad(b, grad, gang)
    elif fill is None: b.fill.background()
    else: b.fill.solid(); b.fill.fore_color.rgb = fill
    if line is None: b.line.fill.background()
    else: b.line.color.rgb = line; b.line.width = Pt(lw)
    return b


def text(s, x, y, w, h, runs, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=1.0, font=FONT):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    if isinstance(runs, str): runs = [(runs, color, bold)]
    p = tf.paragraphs[0]; p.alignment = align; p.line_spacing = space
    for seg in runs:
        t, c, b = seg[0], seg[1], seg[2]; f = seg[3] if len(seg) > 3 else font
        r = p.add_run(); r.text = t; r.font.size = Pt(size); r.font.bold = b; r.font.color.rgb = c; r.font.name = f
    return tb


def bullets(s, x, y, w, h, items, size=13, color=MUT, gap=5, mark=ACC):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.space_after = Pt(gap); p.line_spacing = 1.15
        r0 = p.add_run(); r0.text = "▪  "; r0.font.size = Pt(size); r0.font.color.rgb = mark; r0.font.name = FONT
        r = p.add_run(); r.text = it; r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = FONT
    return tb


def header(s, sec, eyebrow, title_):
    rect(s, ML, Inches(0.5), Inches(0.06), Inches(0.78), fill=GOLD)
    text(s, Emu(ML + Inches(0.2)), Inches(0.48), Inches(9), Inches(0.28), eyebrow, size=11, color=GOLD, bold=True, font=FONT_EN)
    text(s, Emu(ML + Inches(0.2)), Inches(0.76), Inches(10.5), Inches(0.5), title_, size=22, color=INK, bold=True)
    text(s, Inches(11.4), Inches(0.42), Inches(1.5), Inches(0.75), [(sec, GOLD, True, FONT_EN)], size=30, align=PP_ALIGN.RIGHT, font=FONT_EN)
    rect(s, ML, Inches(6.95), CW, Pt(1), fill=LINE)


def footer(s, idx):
    text(s, ML, Inches(7.02), Inches(9.2), Inches(0.28),
         [("上海创智汇 ", SOFT, False), ("· 30场活动 · WAIC2026议程对齐 · 同浦汇", SOFT, False)], size=9)
    rs = text(s, Inches(11.1), Inches(7.02), Inches(1.7), Inches(0.28),
              [("%02d" % idx, GOLD, True, FONT_EN), (" / XX", SOFT, False, FONT_EN)], size=10, align=PP_ALIGN.RIGHT)
    _FOOT.append(rs.text_frame.paragraphs[0].runs[1])


def card(s, x, y, w, h, title=None, items=None, body=None, accent=ACC, tsize=14, bsize=12):
    rect(s, x, y, w, h, grad=[(0, "2A1E55", None), (100, "1E1542", None)], gang=120, line=LINE, lw=1, radius=True)
    rect(s, x, y, Inches(0.07), h, fill=accent)
    if title:
        text(s, Emu(x + Inches(0.25)), Emu(y + Inches(0.14)), Emu(w - Inches(0.4)), Inches(0.38), title, size=tsize, color=INK, bold=True)
    if body:
        text(s, Emu(x + Inches(0.25)), Emu(y + Inches(0.52)), Emu(w - Inches(0.45)), Emu(h - Inches(0.65)), body, size=bsize, color=MUT, space=1.15)
    if items:
        bullets(s, Emu(x + Inches(0.25)), Emu(y + Inches(0.5)), Emu(w - Inches(0.45)), Emu(h - Inches(0.6)), items, size=bsize, mark=accent)


def nxt():
    global IDX; IDX += 1; return IDX


def table_block(s, y0, rows, col_w, sizes=None, rh=0.34):
    yy = y0
    if sizes is None: sizes = [10] * len(rows[0])
    for ri, row in enumerate(rows):
        xx = ML; h = Inches(0.4 if ri == 0 else rh)
        for ci, val in enumerate(row):
            cw = Emu(int(CW * col_w[ci]))
            bg = RGBColor(0x2A, 0x1E, 0x55) if ri == 0 else (RGBColor(0x24, 0x1A, 0x4A) if ri % 2 else RGBColor(0x1C, 0x14, 0x3A))
            rect(s, xx, yy, cw, h, fill=bg, line=LINE, lw=0.45)
            text(s, Emu(xx + Inches(0.05)), Emu(yy + Inches(0.04)), Emu(cw - Inches(0.08)), Emu(h - Inches(0.06)),
                 str(val), size=sizes[ci], color=GOLD if ri == 0 else INK, bold=(ri == 0), anchor=MSO_ANCHOR.MIDDLE)
            xx = Emu(xx + cw)
        yy = Emu(yy + h)
    return yy


# ========== SLIDES ==========
s = slide()
text(s, ML, Inches(1.4), Inches(11), Inches(0.35), [("WAIC 2026 ALIGNED  ·  EVENT PLAN", GOLD, True, FONT_EN)], size=13)
text(s, ML, Inches(1.9), Inches(12), Inches(1.0), "创智汇 · 30场活动方案", size=34, color=INK, bold=True)
text(s, ML, Inches(2.95), Inches(11), Inches(0.4), "结合 WAIC 2026 完整议程升级　｜　智能伙伴 · 共创未来", size=16, color=GOLD, bold=True)
text(s, ML, Inches(3.45), Inches(11), Inches(0.35), "总—分—总　｜　议程映射 · 分场策划 · 招商闭环", size=14, color=MUT)
rect(s, ML, Inches(4.0), Inches(2.4), Pt(3), fill=GOLD)
card(s, ML, Inches(4.35), Inches(3.8), Inches(2.3), "WAIC锚点", accent=GOLD, items=["主题：智能伙伴 共创未来", "时间：7.17–7.20", "三地四馆主会场", "创智汇=杨浦承接会客厅"])
card(s, Inches(4.8), Inches(4.35), Inches(3.8), Inches(2.3), "升级要点", accent=ACC, items=["30场按议程赛道重命名", "7月设WAIC UP!三日联动", "Builders夜+AGI夜落地园区", "会后全年赛道常态化"])
card(s, Inches(8.8), Inches(4.35), Inches(3.7), Inches(2.3), "执行方", accent=GREEN, items=["产业招服：同浦汇", "协办：科技企业联合会", "服务：科技企业服务中心", "学术：复旦相关支持"])

# WAIC议程映射
idx = nxt(); s = slide(); header(s, "01", "WAIC AGENDA MAP", "总 · WAIC 2026议程 → 创智汇活动映射"); footer(s, idx)
rows = [["WAIC/WAICA议程赛道", "创智汇承接方式", "对应场次类型", "招商价值"]]
for t, d in WAIC_META["tracks"]:
    rows.append([t, d[:22], "见分场执行卡", "主题客群精准导入"])
# overwrite with clearer mapping
rows = [
    ["WAIC/WAICA议程赛道", "创智汇承接活动", "建议档期", "招商价值"],
    ["开幕主旨·智能伙伴", "①开年主题沙龙 + ②3 WAIC开放日", "1月 / 7.17", "定调获客"],
    ["Agent智能体/工具评测", "Agent搭建营、营销Agent营", "2月 / 10月", "Agent团队入驻"],
    ["科学多模态智能体 / AI for Science", "多模态工作坊、成果转化日", "3月 / 4月", "科研转化落户"],
    ["算力Infra / GPU集群", "火山引擎×算力实务营、三券沙龙", "3月 / 11月", "高算力企业"],
    ["具身智能空间交互", "具身智能体验日", "4月", "机器人/具身团队"],
    ["青年菁英 / Builders", "春黑客松、Builders Night", "5月 / 7.17晚", "青年创业转化"],
    ["AGI·治理·安全", "AGI之夜、治理沙龙、季度圆桌", "2月/7.18/4月", "高端研究型客户"],
    ["内容创作（错位协同）", "AIGC特训、创作者首发、市集联展", "全年D类", "5F内容去化"],
]
table_block(s, Inches(1.4), rows, [0.28, 0.32, 0.18, 0.22], sizes=[11, 12, 11, 11], rh=0.5)

# 双价值+适配
idx = nxt(); s = slide(); header(s, "01", "DUAL VALUE", "总 · 园区价值 × 招商价值（WAIC化）"); footer(s, idx)
card(s, ML, Inches(1.4), Inches(5.7), Inches(5.25), "给园区的价值", accent=GOLD, items=[
    "把WAIC国家级议题变成园区可感知日常",
    "7月联动抬升政府/媒体可见度",
    "完成载体KPI与活动补贴基础",
    "3F Agent/算力、5F内容双主题坐实",
    "联合会+服务中心+复旦组合壁垒",
], bsize=13)
card(s, Inches(6.85), Inches(1.4), Inches(5.6), Inches(5.25), "给招商的价值", accent=GREEN, items=[
    "大会与学术议程同频，客群更准",
    "开放日+双夜形成72小时转化窗口",
    "场中看场+7日逼单标准化",
    "政策/算力礼包降低决策抗性",
    "年触达900+，漏斗可量化",
], bsize=13)

# 六大主题
idx = nxt(); s = slide(); header(s, "01", "SIX PILLARS", "总 · 六大主题（WAIC对齐后）"); footer(s, idx)
for i, (k, (name, freq, axis, col)) in enumerate(CAT_META.items()):
    x = ML + Inches((i % 3) * 4.0); y = Inches(1.45) + Inches((i // 3) * 2.55)
    card(s, x, y, Inches(3.85), Inches(2.4), f"{k}. {name}　{freq}", accent=col, items=[axis, "详见逐场执行卡"], bsize=13)

# 总表 1-15 / 16-30
for start, end, title_ in ((1, 15, "30场总表（1–15）· WAIC对齐"), (16, 30, "30场总表（16–30）· WAIC对齐")):
    idx = nxt(); s = slide(); header(s, "02", "CALENDAR", title_); footer(s, idx)
    rows = [["序号", "WAIC赛道", "活动名称", "档期", "场地", "招商作用"]]
    for e in EVENTS:
        if start <= e["id"] <= end:
            rows.append([str(e["id"]), e["waic"], e["name"], e["month"], e["floor"], e["lease"][:16]])
    table_block(s, Inches(1.35), rows, [0.07, 0.12, 0.34, 0.14, 0.16, 0.17], sizes=[10, 10, 11, 10, 10, 10], rh=0.33)

# WAIC周专项三页（23/24/25 加细行程）
for eid in (23, 24, 25):
    e = next(x for x in EVENTS if x["id"] == eid)
    idx = nxt(); s = slide(); header(s, "03", "WAIC WEEK", f"分 · WAIC联动精讲｜{e['name']}"); footer(s, idx)
    rect(s, ML, Inches(1.35), CW, Inches(0.65), grad=[(0, "2A1E55", None), (100, "1E1542", None)], line=GOLD, lw=1, radius=True)
    text(s, Emu(ML + Inches(0.25)), Inches(1.42), Inches(11.3), Inches(0.5),
         [(f"{e['month']}　｜　{e['floor']}　｜　{e['dur']}　｜　对齐：{e['waic']}　｜　{e['staff']}", GOLD, True)], size=13, anchor=MSO_ANCHOR.MIDDLE)
    card(s, ML, Inches(2.15), Inches(6.0), Inches(4.5), "具体行程", accent=GOLD, items=e["agenda"], bsize=13)
    card(s, Inches(7.2), Inches(2.15), Inches(5.25), Inches(2.15), "嘉宾与落位", accent=ACC, items=[
        f"嘉宾：{e['guests']}", f"形式：{e['format']}", "落位：签到→主区→看场→洽谈角",
    ], bsize=12)
    card(s, Inches(7.2), Inches(4.45), Inches(5.25), Inches(2.2), "价值链接", accent=GREEN, items=[
        f"园区：{e['park']}", f"招商：{e['lease']}",
    ], bsize=12)

# 主题策划案 6
for k, (name, freq, axis, col) in CAT_META.items():
    evs = [e for e in EVENTS if e["cat"] == k]
    idx = nxt(); s = slide(); header(s, "03", f"THEME {k}", f"分 · 主题策划｜{name}"); footer(s, idx)
    card(s, ML, Inches(1.4), Inches(4.0), Inches(5.25), "定位与WAIC适配", accent=col, items=[
        f"场次：{freq}", f"适配：{axis}", f"代表场：{evs[0]['name']}",
        f"WAIC锚点：{evs[0]['waic']}", "统一：签到→内容→看场→意向→回访",
    ], bsize=12)
    sample = evs[0]
    card(s, Inches(5.2), Inches(1.4), Inches(7.2), Inches(5.25), "标准行程模板", accent=GOLD, items=[
        f"时长：{sample['dur']}｜形式：{sample['format']}", f"嘉宾：{sample['guests']}",
        *[f"行程：{a}" for a in sample["agenda"][:5]], f"人力：{sample['staff']}",
    ], bsize=12)

# 30场执行卡
for e in EVENTS:
    idx = nxt(); s = slide(); header(s, "04", f"E{e['id']:02d}", f"{e['id']:02d}　{e['name']}"); footer(s, idx)
    rect(s, ML, Inches(1.35), CW, Inches(0.65), grad=[(0, "2A1E55", None), (100, "1E1542", None)], line=GOLD, lw=1, radius=True)
    text(s, Emu(ML + Inches(0.2)), Inches(1.42), Inches(11.4), Inches(0.5),
         [(f"WAIC赛道：{e['waic']}　｜　{e['month']}　｜　{e['floor']}　｜　{e['dur']}　｜　{e['format']}　｜　{e['staff']}", GOLD, True)],
         size=12, anchor=MSO_ANCHOR.MIDDLE)
    card(s, ML, Inches(2.15), Inches(4.2), Inches(4.5), "具体行程", accent=GOLD, items=e["agenda"], bsize=12)
    card(s, Inches(4.55), Inches(2.15), Inches(3.9), Inches(2.15), "嘉宾与落位", accent=ACC, items=[
        f"嘉宾：{e['guests']}", f"形式：{e['format']}", f"场地：{e['floor']}", "动线：签到→主区→看场→洽谈",
    ], bsize=11)
    card(s, Inches(8.6), Inches(2.15), Inches(3.95), Inches(2.15), "排期执行", accent=ACC2, items=[
        f"档期：{e['month']}", "T-14嘉宾物料", "T-7邀约确认", "T-1彩排", "T+1~7回访逼单",
    ], bsize=11)
    card(s, Inches(4.55), Inches(4.45), Inches(3.9), Inches(2.2), "园区价值", accent=GREEN, items=[e["park"]], bsize=12)
    card(s, Inches(8.6), Inches(4.45), Inches(3.95), Inches(2.2), "招商价值", accent=ROSE, items=[e["lease"], "看场→意向→报价→签约"], bsize=12)

# 协同逼单转化闭环（压缩3页）
idx = nxt(); s = slide(); header(s, "05", "SUPPORT", "总 · 园区协同支持（人员）"); footer(s, idx)
rows = [
    ["场次类型", "同浦汇", "园区配合", "服务中心/其他", "合计"],
    ["常规训练/沙龙", "2–3", "1–2", "0–2", "4–7"],
    ["政策/治理", "2–3", "1–2", "服务中心2–4", "5–9"],
    ["WAIC开放日", "6", "5 + 物业", "服务中心2+媒体", "15+"],
    ["Builders / AGI夜", "4–5", "2–3", "联合会/复旦2", "8–12"],
    ["嘉年华/漫展级", "6–8", "5–6+安保", "外包", "15+"],
]
table_block(s, Inches(1.45), rows, [0.22, 0.16, 0.2, 0.24, 0.12], sizes=[12, 12, 12, 12, 12], rh=0.55)
text(s, ML, Inches(5.6), CW, Inches(1.0),
     "园区协同清单：场地报备 · 导视停车 · 领导出席 · 意向名单共享 · 会后逼单会议室 · WAIC周期接驳指引。", size=14, color=MUT)

idx = nxt(); s = slide(); header(s, "05", "CLOSE & CONVERT", "总 · 逼单 · 转化 · 招商闭环"); footer(s, idx)
card(s, ML, Inches(1.4), Inches(3.9), Inches(5.25), "逼单", accent=GOLD, items=[
    "每场必含≥20分钟看场", "出口意向表+政策码", "限时礼包当场留资", "WAIC周设绿色通道专台", "T+1~7标准化回访",
], bsize=13)
card(s, Inches(5.1), Inches(1.4), Inches(3.9), Inches(5.25), "转化漏斗", accent=ACC, items=[
    "线索留资率≥60%", "7日触达≥3次", "看场率≥25%", "48小时出报价", "月度签约看板",
], bsize=13)
card(s, Inches(9.2), Inches(1.4), Inches(3.35), Inches(5.25), "闭环", accent=GREEN, items=[
    "活动前：目标场次化", "活动中：建档种草", "活动后：7日逼单", "成交：服务中心接棒", "复盘：反哺下场",
], bsize=13)

idx = nxt(); s = slide(); header(s, "END", "SUMMARY", "结论 · 给对方的一页纸"); footer(s, idx)
card(s, ML, Inches(1.4), Inches(5.7), Inches(5.25), "我们交付", accent=GOLD, items=[
    "30场已按WAIC2026议程赛道重命名",
    "7月WAIC UP!：开放日+Builders夜+AGI夜",
    "逐场行程/嘉宾/落位/价值链接",
    "同浦汇执行 + 逼单转化SOP",
    "全年赛道常态化，不仅办一次大会",
], bsize=13)
card(s, Inches(6.85), Inches(1.4), Inches(5.6), Inches(5.25), "请园区协同", accent=GREEN, items=[
    "确认WAIC周期场地与接驳",
    "按场次配备1–6名现场支持",
    "共享在谈客户名单",
    "开放样板间看场动线",
    "月度联席复盘转化数据",
], bsize=13)

total = len(_FOOT)
for run in _FOOT:
    run.text = " / %02d" % total

os.makedirs(ART, exist_ok=True)
out_cn = os.path.join(HERE, "创智汇30场活动具体方案-WAIC对齐版.pptx")
out_en = os.path.join(HERE, "chuangzhihui-30-events-waic-aligned.pptx")
# also overwrite main deliverable names for convenience
out_main = os.path.join(HERE, "创智汇30场活动具体方案.pptx")
out_main_en = os.path.join(HERE, "chuangzhihui-30-events-plan.pptx")
prs.save(out_cn); prs.save(out_en); prs.save(out_main); prs.save(out_main_en)
for p in (out_cn, out_en, out_main, out_main_en):
    shutil.copy2(p, os.path.join(ART, os.path.basename(p)))
print("PPT slides:", len(prs.slides))

# Excel
wb = Workbook()
ws = wb.active; ws.title = "30场总表-WAIC对齐"
headers = ["序号", "主题", "WAIC赛道", "活动名称", "档期", "场地", "时长", "形式", "嘉宾", "行程", "园区价值", "招商价值", "人力"]
for j, h in enumerate(headers, 1):
    cell = ws.cell(1, j, h); cell.font = Font(name="微软雅黑", bold=True, color="FFFFFF")
    cell.fill = PatternFill("solid", fgColor="4A2C7A"); cell.alignment = Alignment(wrap_text=True, vertical="center")
for i, e in enumerate(EVENTS, 2):
    vals = [e["id"], CAT_META[e["cat"]][0], e["waic"], e["name"], e["month"], e["floor"], e["dur"], e["format"],
            e["guests"], " / ".join(e["agenda"]), e["park"], e["lease"], e["staff"]]
    for j, v in enumerate(vals, 1):
        cell = ws.cell(i, j, v); cell.font = Font(name="微软雅黑", size=9); cell.alignment = Alignment(wrap_text=True, vertical="center")
        if i % 2 == 0: cell.fill = PatternFill("solid", fgColor="F3E9FF")
ws.row_dimensions[1].height = 28
for i in range(2, 32): ws.row_dimensions[i].height = 50
for i, w in enumerate([6, 16, 12, 30, 12, 16, 14, 12, 24, 42, 28, 22, 16], 1):
    ws.column_dimensions[get_column_letter(i)].width = w

ws2 = wb.create_sheet("WAIC议程映射")
for i, row in enumerate([
    ["项目", "内容"],
    ["大会主题", WAIC_META["theme"]],
    ["时间", WAIC_META["when"]],
    ["地点", WAIC_META["where"]],
    ["学术板块", "WAICA 7.18–7.20：主会报告 + 科学多模态智能体 / 数学建模 / 天基智能计算 / 量子加速 / 具身智能 / 青年菁英会"],
    ["夜场逻辑参考", "Builders Night（暖身社交）→ 通往AGI之夜（深度思辨）；创智汇7.17–7.18落地"],
    ["园区角色", "杨浦承接会客厅：大会流量日间承接 + 夜场深度转化 + 全年赛道常态化"],
], 1):
    for j, v in enumerate(row, 1):
        cell = ws2.cell(i, j, v); cell.font = Font(name="微软雅黑", bold=(i == 1 or j == 1), size=10)
        cell.alignment = Alignment(wrap_text=True, vertical="center")
        if i == 1: cell.fill = PatternFill("solid", fgColor="4A2C7A"); cell.font = Font(name="微软雅黑", bold=True, color="FFFFFF")
ws2.column_dimensions["A"].width = 18; ws2.column_dimensions["B"].width = 80
for r in range(2, 8): ws2.row_dimensions[r].height = 36

ws3 = wb.create_sheet("WAICA学术日程摘要")
for i, row in enumerate([
    ["日期", "时段", "议程", "创智汇可借势"],
    ["7.18", "09:00–12:00", "WAICA开幕式", "开放日预热/接驳"],
    ["7.18", "14:00–17:00", "主会技术报告1（最佳论文/Agent等）", "Builders Night议题素材"],
    ["7.19", "09:00–12:00", "科学多模态智能体 / 数学建模与科学计算工作坊", "AI for Science对接日话术"],
    ["7.19", "14:00–17:00", "天基智能计算·智能体 / 量子计算加速", "算力Infra实务营内容"],
    ["7.20", "09:00–12:00", "主会4–6 + 具身智能空间交互 + 青年菁英会", "具身体验日/青年黑客松"],
], 1):
    for j, v in enumerate(row, 1):
        cell = ws3.cell(i, j, v); cell.font = Font(name="微软雅黑", bold=(i == 1), color=("FFFFFF" if i == 1 else "2C2C2C"), size=10)
        if i == 1: cell.fill = PatternFill("solid", fgColor="4A2C7A")
        cell.alignment = Alignment(wrap_text=True, vertical="center")
for i, w in enumerate([10, 14, 42, 28], 1):
    ws3.column_dimensions[get_column_letter(i)].width = w

x_cn = os.path.join(HERE, "创智汇30场活动方案排期表-WAIC对齐.xlsx")
x_en = os.path.join(HERE, "chuangzhihui-30-events-schedule-waic.xlsx")
x_main = os.path.join(HERE, "创智汇30场活动方案排期表.xlsx")
x_main_en = os.path.join(HERE, "chuangzhihui-30-events-schedule.xlsx")
for p in (x_cn, x_en, x_main, x_main_en):
    wb.save(p); shutil.copy2(p, os.path.join(ART, os.path.basename(p)))
print("Excel OK")
