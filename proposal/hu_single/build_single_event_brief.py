# -*- coding: utf-8 -*-
"""胡会长单场活动执行框架 · Word修订 + 独立汇报PPT
口径对齐创智汇定稿：
- 展示「一场」执行与甲乙方分工，档期灵活（非强制一月一场）
- 每场≤30人（修订原稿≥30）
- 甲方=园区（销售/场地）；乙方=同浦汇（策划执行与带客）
- D日建群；回访摘要交园区；园区群内促成成交
"""
import os, shutil
from docx import Document
from docx.shared import Pt, Cm, RGBColor as DocRGB
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Inches, Pt as PPt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn as ppt_qn

HERE = os.path.dirname(os.path.abspath(__file__))
ART = "/opt/cursor/artifacts"
FONT = "Microsoft YaHei"

BG_A, BG_B = "1A1140", "3A206E"
PANEL = RGBColor(0x27, 0x1A, 0x4D)
PANEL2 = RGBColor(0x1E, 0x15, 0x42)
INK = RGBColor(0xF5, 0xF1, 0xFF)
MUT = RGBColor(0xC6, 0xBA, 0xEA)
SOFT = RGBColor(0x94, 0x86, 0xC4)
ACC = RGBColor(0x8B, 0x7B, 0xFF)
GOLD = RGBColor(0xE9, 0xC2, 0x7C)
GREEN = RGBColor(0x56, 0xD6, 0xB6)
LINE = RGBColor(0x4A, 0x3C, 0x7C)
ROW_A = RGBColor(0x2A, 0x1C, 0x52)
ROW_B = RGBColor(0x23, 0x17, 0x48)

SW, SH = Inches(13.333), Inches(7.5)
ML, CW = Inches(0.7), Inches(11.9)

# 单场执行节点（修订版）
NODES = [
    dict(t="T-14（活动前约2周）", act="锁定主题、议程、嘉宾与物料",
         a="审定主题与档期；确认场地/样板间；共享在谈名单（如有）",
         b="提出主题方案与议程草案；定向邀约嘉宾；制作报名物料"),
    dict(t="T-7（活动前1周）", act="邀约确认 + 观众名单 + 场地报备",
         a="审核议程终稿；确认场地物业/导视/安保；接收观众名单",
         b="按产业主题客群邀约观众（目标≤30人）；嘉宾再确认；报备场地需求"),
    dict(t="T-3（活动前3天）", act="演讲PPT与现场物料确认",
         a="审核嘉宾演讲PPT要点；确认桌椅/投影/话筒/名牌到位",
         b="收集整理嘉宾PPT；确认打印物料与动线物料无误"),
    dict(t="T-2（活动前2天）", act="建活动群",
         a="指定专人入群；回答地址/停车等园区侧问题",
         b="拉群、发议程与到场须知；名单分层标注"),
    dict(t="T-1（活动前1天）", act="彩排 + 看场动线走查",
         a="开放样板间/看场动线；物业配合彩排",
         b="全流程彩排；看场带队预案；礼品/政策包就位"),
    dict(t="D日（活动当天）", act="签到→主区→看场→建群",
         a="现场项目介绍/销售支持；开放看场；可在群内发布条件",
         b="接待主持控场；签到分层；带队看场；当场建群与意向登记"),
    dict(t="T+1～7（活动后一周）", act="带客回访 + 销售促成",
         a="基于名单/群做销售跟进与成交交割（园区主责销售）",
         b="主责带客回访；交回访摘要（知会版）；签到表/录音文字稿/AI大纲等资料归档"),
    dict(t="活动后宣发", act="资料宣发与复盘",
         a="公众号/小红书等官方渠道宣发",
         b="提供通稿素材与照片；配合宣发；提交单场复盘纪要"),
]


def build_docx():
    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "微软雅黑"
    style._element.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")
    style.font.size = Pt(11)

    def add_title(text, size=16, bold=True):
        p = doc.add_paragraph()
        r = p.add_run(text)
        r.bold = bold
        r.font.size = Pt(size)
        r.font.name = "微软雅黑"
        r._element.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")
        return p

    add_title("创智汇 · 单场活动执行与甲乙方分工（汇报框架）", 18)
    p = doc.add_paragraph()
    r = p.add_run(
        "用途：内部汇报「一场」活动怎么执行、彼此如何分工。"
        "档期灵活，不强制一个月一场；本框架仅展示单场执行内容。"
        "人数口径：每场到场≤30人。"
    )
    r.font.size = Pt(10.5)
    r.font.name = "微软雅黑"
    r._element.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")

    add_title("一、角色定义", 14)
    for line in [
        "甲方（园区）：方案/档期审定、场地物业与样板间、销售促成与合同交割。",
        "乙方（同浦汇）：策划执行、定向邀约与带客、主持控场、建群、回访摘要。",
        "一句话：我们负责带客，园区负责销售部分；活动结束建群，园区可在群内促成成交。",
        "协同支持（按需）：云创基地、科企联、科技企业服务中心、复旦住房政策研究中心（学术）。",
    ]:
        doc.add_paragraph(line, style=None)

    add_title("二、单场执行流程及甲乙方分工", 14)
    table = doc.add_table(rows=1 + len(NODES), cols=4)
    table.style = "Table Grid"
    headers = ["时间节点", "动作", "甲方（园区）", "乙方（同浦汇）"]
    for i, h in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = h
        for p in cell.paragraphs:
            for r in p.runs:
                r.bold = True
                r.font.size = Pt(10)
                r.font.name = "微软雅黑"
                r._element.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")
    for i, n in enumerate(NODES):
        row = table.rows[i + 1].cells
        vals = [n["t"], n["act"], n["a"], n["b"]]
        for j, v in enumerate(vals):
            row[j].text = v
            for p in row[j].paragraphs:
                for r in p.runs:
                    r.font.size = Pt(9)
                    r.font.name = "微软雅黑"
                    r._element.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")

    add_title("三、单场关键交付物（乙方为主归档）", 14)
    for line in [
        "活动前：主题方案、议程单、嘉宾确认、报名页/物料。",
        "活动中：签到表、意向登记、活动群、看场动线记录。",
        "活动后7天内：回访摘要（知会园区）、照片/通稿素材、录音文字稿与AI大纲（如有）、单场复盘纪要。",
        "说明：交园区的回访材料为摘要知会，便于销售跟进，不额外增加园区事务负担。",
    ]:
        doc.add_paragraph(line)

    add_title("四、口径说明（与总册一致）", 14)
    for line in [
        "人数：每场≤30人；负责人占比适度承兑约50%（非100%）。",
        "客群：以园区产业主题定向邀约为主；弱化关联度核验表表述。",
        "出海/领事：另议另计价，不纳入常规单场年度包默认动作。",
        "不承诺必带外资企业。",
    ]:
        doc.add_paragraph(line)

    add_title("五、单场信息填写栏（汇报时按场填写）", 14)
    fill = doc.add_table(rows=8, cols=2)
    fill.style = "Table Grid"
    fields = [
        ("活动名称", ""),
        ("建议档期", "____年__月__日（灵活，非强制月度）"),
        ("场地", "创智汇 3F / 5F / 其他：______"),
        ("预计人数", "≤30人"),
        ("主题一句话", ""),
        ("目标客群", ""),
        ("本场招商落点", "看场 / 展位 / 入孵 / 其他"),
        ("本场资源导入（如有）", "WAIC切片 / ChinaJoy切片 / 园区自建"),
    ]
    for i, (k, v) in enumerate(fields):
        fill.rows[i].cells[0].text = k
        fill.rows[i].cells[1].text = v
        for c in fill.rows[i].cells:
            for p in c.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(10)
                    r.font.name = "微软雅黑"
                    r._element.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")

    p = doc.add_paragraph()
    r = p.add_run("\n编制参考：同浦汇 × 创智汇活动专项口径　｜　呈阅：胡会长")
    r.font.size = Pt(9)
    r.font.color.rgb = DocRGB(0x66, 0x66, 0x66)

    outs = [
        "创智汇-单场活动执行分工-胡会长汇报框架.docx",
        "chuangzhihui-single-event-ops-framework.docx",
    ]
    os.makedirs(ART, exist_ok=True)
    for fn in outs:
        path = os.path.join(HERE, fn)
        doc.save(path)
        shutil.copy2(path, os.path.join(ART, fn))
        # also parent proposal folder
        parent = os.path.dirname(HERE)
        shutil.copy2(path, os.path.join(parent, fn))
    print("Word OK")


def _grad(shape, stops, ang=90):
    spPr = shape._element.spPr
    for t in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill"):
        e = spPr.find(ppt_qn(t))
        if e is not None:
            spPr.remove(e)
    g = spPr.makeelement(ppt_qn("a:gradFill"), {})
    lst = g.makeelement(ppt_qn("a:gsLst"), {})
    for pos, col, al in stops:
        gs = g.makeelement(ppt_qn("a:gs"), {"pos": str(int(pos * 1000))})
        c = g.makeelement(ppt_qn("a:srgbClr"), {"val": col})
        if al is not None:
            c.append(g.makeelement(ppt_qn("a:alpha"), {"val": str(int(al * 1000))}))
        gs.append(c); lst.append(gs)
    g.append(lst)
    g.append(g.makeelement(ppt_qn("a:lin"), {"ang": str(int(ang * 60000)), "scaled": "1"}))
    ln = spPr.find(ppt_qn("a:ln"))
    (ln.addprevious(g) if ln is not None else spPr.append(g))


def slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.shadow.inherit = False; bg.line.fill.background()
    _grad(bg, [(0, BG_A, None), (55, BG_B, None), (100, BG_A, None)], 120)
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s


def rect(s, x, y, w, h, fill=None, line=None, lw=1.0, radius=False):
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    b.shadow.inherit = False
    if fill is None: b.fill.background()
    else: b.fill.solid(); b.fill.fore_color.rgb = fill
    if line is None: b.line.fill.background()
    else: b.line.color.rgb = line; b.line.width = PPt(lw)
    return b


def text(s, x, y, w, h, content, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = PPt(2); tf.margin_right = PPt(2); tf.margin_top = PPt(1); tf.margin_bottom = PPt(1)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = content; r.font.size = PPt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = FONT
    return tb


def bullets(s, x, y, w, h, items, size=13, color=MUT, mark=ACC):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = PPt(5); p.line_spacing = 1.15
        r0 = p.add_run(); r0.text = "▸ "; r0.font.size = PPt(size); r0.font.color.rgb = mark; r0.font.name = FONT
        r = p.add_run(); r.text = it; r.font.size = PPt(size); r.font.color.rgb = color; r.font.name = FONT
    return tb


def header(s, sec, title, sub=None):
    text(s, ML, Inches(0.28), Inches(12), Inches(0.28), sec, size=13, color=ACC, bold=True)
    text(s, ML, Inches(0.55), Inches(12), Inches(0.42), title, size=24, color=INK, bold=True)
    rect(s, ML, Inches(1.05), Inches(0.55), Inches(0.04), fill=GOLD)
    if sub:
        text(s, ML, Inches(1.18), Inches(12), Inches(0.28), sub, size=12, color=MUT)


def footer(s, page, total):
    text(s, ML, Inches(7.1), Inches(9.5), Inches(0.25),
         "创智汇 × 同浦汇 · 单场活动执行汇报框架 · 呈阅胡会长", size=11, color=SOFT)
    text(s, Inches(11.1), Inches(7.1), Inches(1.7), Inches(0.25),
         f"{page:02d} / {total:02d}", size=11, color=SOFT, align=PP_ALIGN.RIGHT)


def card(s, x, y, w, h, title, items, accent=ACC):
    rect(s, x, y, w, h, fill=PANEL, line=LINE, radius=True)
    rect(s, x, y, Inches(0.08), h, fill=accent)
    text(s, Emu(x + Inches(0.25)), Emu(y + Inches(0.15)), Emu(w - Inches(0.4)), Inches(0.32),
         title, size=15, color=INK, bold=True)
    bullets(s, Emu(x + Inches(0.25)), Emu(y + Inches(0.55)), Emu(w - Inches(0.45)), Emu(h - Inches(0.7)),
            items, size=13, mark=accent)


def table(s, y0, rows, col_w, sizes=None, rh=0.34):
    yy = y0
    if sizes is None: sizes = [11] * len(rows[0])
    for ri, row in enumerate(rows):
        xx = ML; h = Inches(0.4 if ri == 0 else rh)
        for ci, val in enumerate(row):
            cw = Emu(int(CW * col_w[ci]))
            if ri == 0:
                bg, fg, bd = PANEL, GOLD, True
            else:
                bg = ROW_A if ri % 2 == 0 else ROW_B
                fg, bd = INK, False
            rect(s, xx, yy, cw, h, fill=bg, line=LINE, lw=0.5)
            text(s, Emu(xx + Inches(0.05)), Emu(yy + Inches(0.04)), Emu(cw - Inches(0.08)), Emu(h - Inches(0.06)),
                 str(val), size=sizes[ci], color=fg, bold=bd, anchor=MSO_ANCHOR.MIDDLE)
            xx = Emu(xx + cw)
        yy = Emu(yy + h)
    return yy


def build_ppt():
    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
    TOTAL = 10
    n = 0

    def pg(s):
        nonlocal n
        n += 1
        footer(s, n, TOTAL)
        return n

    # 1 cover
    s = slide(prs)
    text(s, ML, Inches(1.7), Inches(12), Inches(0.3), "创智汇 × 同浦汇｜内部汇报框架　·　呈阅胡会长", size=14, color=ACC, bold=True)
    text(s, ML, Inches(2.2), Inches(12), Inches(0.55), "单场活动执行与甲乙方分工", size=34, color=INK, bold=True)
    text(s, ML, Inches(2.9), Inches(12), Inches(0.35), "展示「一场」怎么做　｜　档期灵活　｜　彼此分工说清楚", size=16, color=GOLD, bold=True)
    text(s, ML, Inches(3.4), Inches(12), Inches(0.35),
         "每场≤30人　·　同浦汇带客 / 园区销售　·　活动结束建群", size=14, color=MUT)
    card(s, ML, Inches(4.1), Inches(5.8), Inches(2.5), "本册是什么", [
        "不是全年30场总册",
        "只讲「单场」执行流程与分工",
        "可套用到任意一场活动汇报",
        "基于胡会长初版框架修订",
    ], ACC)
    card(s, Inches(7.2), Inches(4.1), Inches(5.5), Inches(2.5), "相对初版的调整", [
        "≥30人 → 改为≤30人",
        "「一月一场」→ 改为灵活档期示例",
        "明确甲方=园区、乙方=同浦汇",
        "补齐建群、回访摘要、销售促成",
    ], GOLD)
    pg(s)

    # 2 purpose
    s = slide(prs); header(s, "01 · 使用说明", "这份框架怎么用", "内部汇报单场活动情况；不一定一个月一场")
    card(s, ML, Inches(1.5), Inches(5.8), Inches(5.2), "适用场景", [
        "向会长/园区汇报某一场怎么执行",
        "确认本场甲乙方谁做什么",
        "活动前对齐节点与交付物",
        "活动后对照复盘是否闭环",
    ], ACC)
    card(s, Inches(7.2), Inches(1.5), Inches(5.5), Inches(5.2), "怎么填", [
        "先定：主题 / 档期 / 场地 / 客群",
        "再套：T-14→D日→T+7 分工表",
        "当场：签到、看场、建群",
        "事后：回访摘要交园区销售跟进",
    ], GREEN)
    pg(s)

    # 3 roles
    s = slide(prs); header(s, "02 · 角色定义", "甲方园区 · 乙方同浦汇 · 一句话分工")
    card(s, ML, Inches(1.5), Inches(5.8), Inches(5.2), "甲方（园区）", [
        "审定主题、议程与档期",
        "场地、物业、导视、样板间",
        "共享在谈名单（如有）",
        "负责销售部分与合同交割",
        "可在活动群内发布条件促成成交",
    ], GOLD)
    card(s, Inches(7.2), Inches(1.5), Inches(5.5), Inches(5.2), "乙方（同浦汇）", [
        "策划主题与议程草案",
        "定向邀约嘉宾与观众（带客）",
        "主持控场、看场带队",
        "建群 + 回访摘要主责",
        "资料归档与复盘纪要",
    ], ACC)
    text(s, ML, Inches(6.85), CW, Inches(0.3),
         "一句话：我们负责带客，园区负责销售部分。协同支持按需：云创基地 / 科企联 / 服中心 / 复旦住房中心。",
         size=13, color=MUT)
    pg(s)

    # 4 timeline overview
    s = slide(prs); header(s, "03 · 单场时间轴", "同一套打法（可套任意一场）")
    rows = [["节点", "关键动作", "结果"]]
    for n0 in NODES:
        rows.append([n0["t"].split("（")[0], n0["act"][:18], "可核验节点"])
    table(s, Inches(1.4), rows, [0.28, 0.42, 0.3], sizes=[12, 13, 12], rh=0.48)
    pg(s)

    # 5-6 division tables
    for subset, title in ((NODES[:4], "单场分工表（上）· 筹备阶段"), (NODES[4:], "单场分工表（下）· 现场与复盘")):
        s = slide(prs); header(s, "04 · 甲乙方分工", title, "在胡会长初版框架上补齐建群/回访/销售促成")
        rows = [["时间", "动作", "甲方（园区）", "乙方（同浦汇）"]]
        for n0 in subset:
            rows.append([n0["t"][:12], n0["act"][:14], n0["a"][:22], n0["b"][:22]])
        table(s, Inches(1.4), rows, [0.18, 0.2, 0.31, 0.31], sizes=[11, 11, 11, 11], rh=0.7)
        pg(s)

    # 7 deliverables
    s = slide(prs); header(s, "05 · 单场交付物", "交给园区的、可核验的东西")
    table(s, Inches(1.45), [
        ["阶段", "交付物", "责任"],
        ["活动前", "主题方案、议程单、嘉宾确认、报名页/物料", "乙方制作；甲方审定"],
        ["活动中", "签到表、意向登记、活动群、看场记录", "乙方执行；甲方现场支持"],
        ["活动后7天", "回访摘要（知会版）、照片/通稿、文字稿/AI大纲、复盘纪要", "乙方归档交园区"],
        ["销售促成", "房源条件、报价、合同交割", "甲方（园区）主责"],
    ], [0.16, 0.54, 0.3], sizes=[13, 13, 13], rh=0.7)
    text(s, ML, Inches(6.5), CW, Inches(0.4),
         "回访摘要为知会版：方便园区销售跟进，不额外增加园区事务负担。", size=13, color=MUT)
    pg(s)

    # 8 fill template
    s = slide(prs); header(s, "06 · 本场信息栏", "汇报时按场填写（模板）")
    table(s, Inches(1.45), [
        ["字段", "填写示例 / 留白"],
        ["活动名称", "例如：创智汇项目推广日①·中介与渠道专场"],
        ["档期", "____年__月__日（灵活，非强制月度）"],
        ["场地", "创智汇3F / 5F / 其他______"],
        ["预计人数", "≤30人　｜　负责人约50%"],
        ["主题一句话", "________________"],
        ["目标客群", "________________"],
        ["招商落点", "看场 / 展位 / 入孵 / 其他"],
        ["资源导入（如有）", "WAIC切片 / ChinaJoy切片 / 园区自建"],
    ], [0.28, 0.72], sizes=[13, 13], rh=0.48)
    pg(s)

    # 9口径
    s = slide(prs); header(s, "07 · 口径对齐", "与创智汇活动总册保持一致")
    card(s, ML, Inches(1.5), Inches(5.8), Inches(5.2), "人数与客群", [
        "每场到场≤30人",
        "负责人占比适度承兑约50%",
        "产业主题定向邀约为主",
        "弱化关联度核验表表述",
        "不承诺必带外资企业",
    ], ACC)
    card(s, Inches(7.2), Inches(1.5), Inches(5.5), Inches(5.2), "边界", [
        "出海活动：另议另计价",
        "领事到访/挂牌：另计价且分性质",
        "本框架不含租金对赌条款",
        "本框架不是全年排期总表",
        "只服务「单场执行汇报」",
    ], GOLD)
    pg(s)

    # 10 end
    s = slide(prs)
    text(s, ML, Inches(2.4), Inches(12), Inches(0.5), "谢谢审阅", size=34, color=INK, bold=True, align=PP_ALIGN.CENTER)
    text(s, ML, Inches(3.2), Inches(12), Inches(0.4), "单场活动执行与甲乙方分工 · 汇报框架", size=18, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    text(s, ML, Inches(3.8), Inches(12), Inches(0.35), "同浦汇带客 · 园区销售 · 一场一套打法", size=14, color=MUT, align=PP_ALIGN.CENTER)
    text(s, ML, Inches(4.5), Inches(12), Inches(0.8),
         "呈阅：胡会长\n编制参考：同浦汇（基于初版框架修订）",
         size=13, color=SOFT, align=PP_ALIGN.CENTER)
    pg(s)

    assert n == TOTAL
    outs = [
        "创智汇-单场活动执行分工-胡会长汇报.pptx",
        "chuangzhihui-single-event-ops-brief.pptx",
    ]
    tmp = os.path.join(HERE, "_tmp_single.pptx")
    prs.save(tmp)
    parent = os.path.dirname(HERE)
    os.makedirs(ART, exist_ok=True)
    for fn in outs:
        shutil.copy2(tmp, os.path.join(HERE, fn))
        shutil.copy2(tmp, os.path.join(parent, fn))
        shutil.copy2(tmp, os.path.join(ART, fn))
    os.remove(tmp)
    print("PPT slides:", len(prs.slides))


if __name__ == "__main__":
    build_docx()
    build_ppt()
