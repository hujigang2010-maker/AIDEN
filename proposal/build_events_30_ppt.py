# -*- coding: utf-8 -*-
"""创智汇 · 30场活动具体方案 PPT（总—分—总）
含：主题命名、园区/招商价值、主题适配、分场策划（行程/嘉宾/落位/排期/价值）、
协同支持（人员/逼单/转化/招商闭环）
用法: python3 build_events_30_ppt.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

HERE = os.path.dirname(os.path.abspath(__file__))
ART = "/opt/cursor/artifacts"
FONT, FONT_EN = "Microsoft YaHei", "Arial"
BG_A, BG_B = "1A1140", "3A206E"
INK = RGBColor(0xF5, 0xF1, 0xFF)
MUT = RGBColor(0xC6, 0xBA, 0xEA)
SOFT = RGBColor(0x94, 0x86, 0xC4)
ACC = RGBColor(0x8B, 0x7B, 0xFF)
ACC2 = RGBColor(0xC7, 0x7D, 0xFF)
GOLD = RGBColor(0xE9, 0xC2, 0x7C)
GREEN = RGBColor(0x56, 0xD6, 0xB6)
ROSE = RGBColor(0xF5, 0x8B, 0xB8)
LINE = RGBColor(0x4A, 0x3C, 0x7C)
SW, SH = Inches(13.333), Inches(7.5)
ML, CW = Inches(0.75), Inches(11.83)
prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]
_FOOT = []
IDX = 0

# —— 30 场活动数据 ——
# cat: A培训 B政策 C对接 D潮玩 E大型 F出海
EVENTS = [
    # A ×10 AI培训/黑客松
    dict(id=1, cat="A", name="创智汇·AI工具链上手营", month="1月下旬", floor="3F培训/OPC",
         dur="半天（13:30–17:30）", guests="AI讲师×1、园区政策专员×1",
         agenda=["13:30 签到参观样板间", "14:00 开场+空间政策礼包", "14:30 工具链实操", "16:30 工位/单元看场", "17:00 意向登记逼单"],
         park="OPC社群体验、载体活跃度、培训KPI", lease="引流AI初创，当场看场转化",
         format="工作坊+看场", staff="我方3+园区2"),
    dict(id=2, cat="A", name="AIGC内容生产实战工作坊", month="2月中旬", floor="3F直播间+培训",
         dur="1天", guests="AIGC创作者、内容机构导师",
         agenda=["09:30 签到", "10:00 案例拆解", "13:30 实操产出", "16:00 优秀作品路演", "16:40 入驻政策+看场"],
         park="内容氛围、直播间利用率", lease="锁定内容型租户与联合办公需求",
         format="实训营", staff="我方4+园区2"),
    dict(id=3, cat="A", name="OPC超级个体黑客松（春季）", month="3月上旬", floor="3F整层开放",
         dur="1.5天", guests="技术评委×3、投资人×2、云厂商",
         agenda=["D1 组队开题", "D1夜 封闭开发", "D2路演决赛", "颁奖+入驻礼包发布", "一对一看场"],
         park="品牌声量、媒体曝光、青年人才黏性", lease="获奖团队优先谈单元/工位",
         format="黑客松", staff="我方5+园区3"),
    dict(id=4, cat="A", name="数字人直播实训营", month="3月下旬", floor="3F直播间",
         dur="半天", guests="数字人厂商、直播操盘手",
         agenda=["参观直播间", "数字人搭建演示", "跟播实操", "商业变现路径", "空间与套餐报价"],
         park="直播间运营样板", lease="直播/MCN类客户转化",
         format="实训+演示", staff="我方3+园区2"),
    dict(id=5, cat="A", name="火山引擎·算力应用训练营", month="4月中旬", floor="3F+云创协同",
         dur="1天", guests="火山引擎讲师、算力政策专员",
         agenda=["算力政策包解读", "代金券/折扣实操", "企业案例", "一对一算力诊断", "入驻+云资源捆绑报价"],
         park="专有算力政策兑现、云基地协同", lease="AI算力消耗型企业定向招商",
         format="厂商联训", staff="我方3+园区2+厂商2"),
    dict(id=6, cat="A", name="AI微短剧脚本与制作特训", month="5月上旬", floor="5F展厅角+3F",
         dur="1天", guests="微短剧导演/编剧、平台方",
         agenda=["政策与集聚区解读", "脚本工作坊", "AIGC制片流程", "优秀剧本路演", "5F场景踩点"],
         park="对齐YOUNG立方/微短剧定位", lease="厂牌/工作室看场5F或3F",
         format="特训营", staff="我方4+园区2"),
    dict(id=7, cat="A", name="智能体搭建一日营", month="6月中旬", floor="3F OPC",
         dur="1天", guests="智能体工程师、行业顾问",
         agenda=["需求拆解", "低代码搭建", "行业垂类案例", "作品评审", "入驻+服务包"],
         park="OPC超级个体社区氛围", lease="工具型/服务型小团队入驻",
         format="一日营", staff="我方3+园区2"),
    dict(id=8, cat="A", name="OPC超级个体黑客松（秋季）", month="9月中旬", floor="3F整层",
         dur="1.5天", guests="评委+赞助商+媒体",
         agenda=["同春季结构", "增加秋季赛道（出海/短剧）", "决赛路演", "媒体采访", "集中逼单日"],
         park="年度品牌IP化第二峰", lease="秋季去化冲刺客源",
         format="黑客松", staff="我方5+园区3"),
    dict(id=9, cat="A", name="AI营销投放实战营", month="10月中旬", floor="3F培训",
         dur="半天", guests="投放操盘手、品牌主",
         agenda=["投放策略", "素材AIGC生产", "投流复盘", "园区企业案例", "增值服务转化"],
         park="企业服务增收、会员粘性", lease="营销型公司与广告客户",
         format="实战营", staff="我方3+园区1"),
    dict(id=10, cat="A", name="创智汇AI年度Demo Day", month="12月上旬", floor="5F主展区",
         dur="1天", guests="投资人、链主、媒体、区级嘉宾",
         agenda=["开幕致辞", "十强路演", "颁奖", "招商政策发布", "红酒洽谈+看场"],
         park="年度成果展、政府可见度", lease="集中签约意向、媒体背书",
         format="路演日", staff="我方6+园区4"),
    # B ×6 政策
    dict(id=11, cat="B", name="YOUNG立方政策解读沙龙", month="1月中旬", floor="5F沙龙区",
         dur="2.5小时", guests="区政策宣讲人、服务中心",
         agenda=["政策要点", "适用企业画像", "申报路径", "问答", "诊断预约+看场"],
         park="政策礼包前置、载体认知", lease="内容/微短剧企业定向导入",
         format="沙龙", staff="我方2+园区2+服务中心2"),
    dict(id=12, cat="B", name="高企认定冲刺辅导会", month="2月下旬", floor="3F会议室",
         dur="半天", guests="高企辅导老师、财税顾问",
         agenda=["条件拆解", "材料清单", "失败案例", "一对一初筛", "入驻即享辅导套餐"],
         park="在园企业资质提升", lease="外区待认定企业「带政策入驻」",
         format="辅导会", staff="我方2+服务中心3"),
    dict(id=13, cat="B", name="专精特新培育路演", month="4月上旬", floor="5F",
         dur="半天", guests="专精特新评审顾问、银行",
         agenda=["梯度政策", "企业路演", "导师点评", "金融对接", "看场与意向表"],
         park="优质企业筛选与培育", lease="成长型企业补位招商",
         format="路演", staff="我方3+园区2"),
    dict(id=14, cat="B", name="杨浦科创政策一对一诊断日", month="5月下旬", floor="3F多间工位",
         dur="1天滚动", guests="服务中心顾问团",
         agenda=["预约排号", "20分钟/企诊断", "政策匹配表输出", "当场报价", "三日回访逼单"],
         park="服务中心专业形象", lease="高意向客户集中转化",
         format="门诊式诊断", staff="我方3+服务中心4+园区1"),
    dict(id=15, cat="B", name="投融资对接路演（AI赛道）", month="7月中旬", floor="5F",
         dur="半天", guests="基金×5、明星项目×8",
         agenda=["项目路演", "闭门对接", "园区资本资源介绍", "入驻条件", "投资+租赁联动"],
         park="资本氛围、优质项目落位", lease="获投项目优先谈租金/免租",
         format="路演+闭门", staff="我方4+园区2"),
    dict(id=16, cat="B", name="创新券·三券申用实务沙龙", month="11月上旬", floor="3F",
         dur="2小时", guests="创新券平台、云厂商",
         agenda=["券种规则", "核销实操", "案例", "服务机构对接", "入驻企业开户引导"],
         park="创新券流水、服务机构KPI", lease="需用券企业向园区聚集",
         format="实务沙龙", staff="我方2+服务中心2"),
    # C ×6 对接
    dict(id=17, cat="C", name="汕头玩具供应链对接会", month="3月中旬", floor="5F玩具集群区",
         dur="半天", guests="汕头商协会、品牌采购",
         agenda=["集群介绍", "供需对接", "展位参观", "意向签约墙", "展位/仓储报价"],
         park="5F产业集群填充", lease="玩具企业展位去化",
         format="对接会", staff="我方3+园区2"),
    dict(id=18, cat="C", name="扬州毛绒产业游学对接", month="4月下旬", floor="5F+外部游学",
         dur="1天（含半天园内）", guests="扬州产业集群代表",
         agenda=["园内展区介绍", "合作模式（租金/扣点）", "游学小结", "入驻MOU", "跟进清单"],
         park="毛绒集群落地", lease="700㎡量级客户",
         format="游学对接", staff="我方3+园区2"),
    dict(id=19, cat="C", name="东莞潮玩品牌入沪推介", month="6月上旬", floor="5F潮玩区",
         dur="半天", guests="潮玩品牌方、渠道商",
         agenda=["上海渠道机会", "5F落位方案", "联名活动设想", "看场", "定金/意向"],
         park="潮玩内容填充", lease="品牌展位与快闪",
         format="推介会", staff="我方3+园区2"),
    dict(id=20, cat="C", name="IP授权交易撮合会", month="7月下旬", floor="5F展示中心",
         dur="1天", guests="IP方、被授权商、律师",
         agenda=["IP路演", "一对一撮合桌", "法律速递", "成交看板", "落户激励"],
         park="IP交易氛围、分成收入", lease="IP公司与衍生品团队入驻",
         format="撮合会", staff="我方4+园区2"),
    dict(id=21, cat="C", name="高校成果转化对接日", month="9月下旬", floor="3F+沙龙",
         dur="半天", guests="复旦/同济技术转移、教授团队",
         agenda=["成果路演", "中试/场景需求", "园区承接方案", "看孵化单元", "转化落位协议"],
         park="高校协同、学术支持落地", lease="成果公司/实验室落户",
         format="对接日", staff="我方3+园区2+高校2"),
    dict(id=22, cat="C", name="跨境展销买家对接会", month="11月中旬", floor="5F展贸",
         dur="半天", guests="跨境买家、货代、支付机构",
         agenda=["买家需求发布", "供应商闪见", "展位下单", "物流金融配套", "展贸租户转化"],
         park="5F贸易属性兑现", lease="外贸/跨境团队办公+展位",
         format="买家对接", staff="我方4+园区2"),
    # D ×4 IP潮玩
    dict(id=23, cat="D", name="IP首发发布会①", month="5月中旬", floor="5F主展",
         dur="3小时", guests="IP主理人、媒体、渠道",
         agenda=["揭幕", "发布会", "签售/快闪", "媒体专访", "招商通道开放"],
         park="内容事件、客流", lease="关联品牌问询与展位",
         format="发布会", staff="我方4+园区3"),
    dict(id=24, cat="D", name="潮玩主题周末市集", month="6月下旬", floor="5F+公区",
         dur="2天", guests="摊主×40、达人",
         agenda=["摊位布置", "开市", "互动打卡", "夜间场", "优质摊主转正谈"],
         park="人气、经营性收入", lease="摊主升级为固定展位/办公",
         format="市集", staff="我方5+园区4"),
    dict(id=25, cat="D", name="沉浸式IP联展", month="8月中旬", floor="5F展示中心",
         dur="7–10天巡展", guests="联合IP×5、文旅渠道",
         agenda=["布展", "开幕礼", "预约观展", "教育场", "闭幕招商酒会"],
         park="长期人流、品牌形象", lease="闭幕酒会集中转化",
         format="联展", staff="我方4+园区3+执行外包"),
    dict(id=26, cat="D", name="IP首发发布会②", month="10月下旬", floor="5F",
         dur="3小时", guests="IP方、玩具渠道、媒体",
         agenda=["同①结构", "叠加双十一预售", "渠道对接", "看场", "年框意向"],
         park="下半年内容高峰", lease="旺季补位招商",
         format="发布会", staff="我方4+园区3"),
    # E ×2 大型
    dict(id=27, cat="E", name="创智汇春季漫展", month="4月下旬", floor="5F+可临建",
         dur="2天", guests="Coser、品牌商、媒体、赞助商",
         agenda=["开幕", "舞台演出", "展商区", "门票核销", "闭幕+企业招商通道"],
         park="千人级声量、门票赞助收入、载体KPI", lease="展商/品牌方留资转租赁",
         format="漫展", staff="我方8+园区6+安保票务外包"),
    dict(id=28, cat="E", name="创智汇秋季潮玩嘉年华", month="10月上旬", floor="5F+公区",
         dur="2天", guests="潮玩品牌、达人、赞助商",
         agenda=["开幕秀", "市集+展洽", "赛事", "夜间场", "嘉年华签约仪式"],
         park="秋季品牌大事件", lease="集中签约展位/办公",
         format="嘉年华", staff="我方8+园区6+外包"),
    # F ×2 出海
    dict(id=29, cat="F", name="国别会客厅·领事专题沙龙", month="5月下旬", floor="会客厅/5F",
         dur="2.5小时", guests="领事官员、涉外企业、翻译",
         agenda=["外事礼仪接待", "国别机会分享", "B2B闪见", "园区国际服务介绍", "高端看场"],
         park="高端背书、国际形象", lease="涉外/出海企业入驻",
         format="领事沙龙", staff="我方4+园区2+外事协同"),
    dict(id=30, cat="F", name="企业出海推介与买家团", month="11月下旬", floor="5F+培训",
         dur="1天", guests="海外买家、跨境平台、出海顾问",
         agenda=["出海路径课", "买家需求", "一对一洽谈", "订单意向", "办公/展位捆绑"],
         park="出海服务闭环、撮合收入", lease="有出海需求的制造/品牌入驻",
         format="推介+买家团", staff="我方5+园区2"),
]

CAT_META = {
    "A": ("AI培训 / 黑客松", "10场", "3F AI/OPC主轴", GOLD, "引流→实操→看场→工位/单元"),
    "B": ("政策沙龙 / 路演", "6场", "政策礼包主轴", ACC, "解读→诊断→申报→带政策入驻"),
    "C": ("行业对接 / 撮合", "6场", "5F集群+高校", ACC2, "供需匹配→展位/落位→订单或租赁"),
    "D": ("IP / 潮玩活动", "4场", "5F内容主轴", GREEN, "内容事件→客流→展位/品牌入驻"),
    "E": ("漫展 / 大型活动", "2场", "声量与经营", ROSE, "大流量→留资→集中逼单签约"),
    "F": ("领事 / 出海", "2场", "国际与出海", GOLD, "高端背书→涉外企业→出海服务包"),
}


def _grad(shape, stops, ang=90):
    spPr = shape._element.spPr
    for t in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill"):
        e = spPr.find(qn(t))
        if e is not None:
            spPr.remove(e)
    g = spPr.makeelement(qn("a:gradFill"), {})
    lst = g.makeelement(qn("a:gsLst"), {})
    for pos, col, al in stops:
        gs = g.makeelement(qn("a:gs"), {"pos": str(int(pos * 1000))})
        c = g.makeelement(qn("a:srgbClr"), {"val": col})
        if al is not None:
            c.append(g.makeelement(qn("a:alpha"), {"val": str(int(al * 1000))}))
        gs.append(c)
        lst.append(gs)
    g.append(lst)
    g.append(g.makeelement(qn("a:lin"), {"ang": str(int(ang * 60000)), "scaled": "1"}))
    ln = spPr.find(qn("a:ln"))
    (ln.addprevious(g) if ln is not None else spPr.append(g))


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.shadow.inherit = False
    bg.line.fill.background()
    _grad(bg, [(0, BG_A, None), (55, BG_B, None), (100, BG_A, None)], 120)
    s.shapes._spTree.remove(bg._element)
    s.shapes._spTree.insert(2, bg._element)
    return s


def rect(s, x, y, w, h, fill=None, line=None, lw=1.0, radius=False, grad=None, gang=90):
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    b.shadow.inherit = False
    if grad is not None:
        _grad(b, grad, gang)
    elif fill is None:
        b.fill.background()
    else:
        b.fill.solid()
        b.fill.fore_color.rgb = fill
    if line is None:
        b.line.fill.background()
    else:
        b.line.color.rgb = line
        b.line.width = Pt(lw)
    return b


def text(s, x, y, w, h, runs, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, space=1.0, font=FONT):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    if isinstance(runs, str):
        runs = [(runs, color, bold)]
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = space
    for seg in runs:
        t, c, b = seg[0], seg[1], seg[2]
        f = seg[3] if len(seg) > 3 else font
        r = p.add_run()
        r.text = t
        r.font.size = Pt(size)
        r.font.bold = b
        r.font.color.rgb = c
        r.font.name = f
    return tb


def bullets(s, x, y, w, h, items, size=13, color=MUT, gap=5, mark=ACC):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.15
        r0 = p.add_run()
        r0.text = "▪  "
        r0.font.size = Pt(size)
        r0.font.color.rgb = mark
        r0.font.name = FONT
        r = p.add_run()
        r.text = it
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = FONT
    return tb


def header(s, sec, eyebrow, title):
    rect(s, ML, Inches(0.5), Inches(0.06), Inches(0.78), fill=GOLD)
    text(s, Emu(ML + Inches(0.2)), Inches(0.48), Inches(9), Inches(0.28),
         eyebrow, size=11, color=GOLD, bold=True, font=FONT_EN)
    text(s, Emu(ML + Inches(0.2)), Inches(0.76), Inches(10.5), Inches(0.5),
         title, size=22, color=INK, bold=True)
    text(s, Inches(11.4), Inches(0.42), Inches(1.5), Inches(0.75),
         [(sec, GOLD, True, FONT_EN)], size=32, align=PP_ALIGN.RIGHT, font=FONT_EN)
    rect(s, ML, Inches(6.95), CW, Pt(1), fill=LINE)


def footer(s, idx):
    text(s, ML, Inches(7.02), Inches(9), Inches(0.28),
         [("上海创智汇 ", SOFT, False), ("· 30场活动方案 · 同浦汇产业招服", SOFT, False)], size=9)
    rs = text(s, Inches(11.1), Inches(7.02), Inches(1.7), Inches(0.28),
              [("%02d" % idx, GOLD, True, FONT_EN), (" / XX", SOFT, False, FONT_EN)],
              size=10, align=PP_ALIGN.RIGHT)
    _FOOT.append(rs.text_frame.paragraphs[0].runs[1])


def card(s, x, y, w, h, title=None, items=None, body=None, accent=ACC, tsize=14, bsize=12):
    rect(s, x, y, w, h, grad=[(0, "2A1E55", None), (100, "1E1542", None)], gang=120, line=LINE, lw=1, radius=True)
    rect(s, x, y, Inches(0.07), h, fill=accent)
    if title:
        text(s, Emu(x + Inches(0.25)), Emu(y + Inches(0.14)), Emu(w - Inches(0.4)), Inches(0.38),
             title, size=tsize, color=INK, bold=True)
    if body:
        text(s, Emu(x + Inches(0.25)), Emu(y + Inches(0.52)), Emu(w - Inches(0.45)), Emu(h - Inches(0.65)),
             body, size=bsize, color=MUT, space=1.15)
    if items:
        bullets(s, Emu(x + Inches(0.25)), Emu(y + Inches(0.5)), Emu(w - Inches(0.45)), Emu(h - Inches(0.6)),
                items, size=bsize, mark=accent)


def nxt():
    global IDX
    IDX += 1
    return IDX


def add_table_slide(title_zh, rows, col_w, sizes=None, rh=0.38):
    idx = nxt()
    s = slide()
    header(s, "02", "CALENDAR", title_zh)
    footer(s, idx)
    yy = Inches(1.4)
    n = len(rows[0])
    if sizes is None:
        sizes = [10] * n
    for ri, row in enumerate(rows):
        xx = ML
        h = Inches(0.42 if ri == 0 else rh)
        for ci, val in enumerate(row):
            cw = Emu(int(CW * col_w[ci]))
            bg = RGBColor(0x2A, 0x1E, 0x55) if ri == 0 else (
                RGBColor(0x24, 0x1A, 0x4A) if ri % 2 else RGBColor(0x1C, 0x14, 0x3A))
            rect(s, xx, yy, cw, h, fill=bg, line=LINE, lw=0.5)
            text(s, Emu(xx + Inches(0.06)), Emu(yy + Inches(0.05)), Emu(cw - Inches(0.1)), Emu(h - Inches(0.08)),
                 str(val), size=sizes[ci], color=GOLD if ri == 0 else INK,
                 bold=(ri == 0 or ci <= 1), anchor=MSO_ANCHOR.MIDDLE)
            xx = Emu(xx + cw)
        yy = Emu(yy + h)
    return s


# ===================== BUILD SLIDES =====================

# 01 Cover
s = slide()
text(s, ML, Inches(1.5), Inches(11), Inches(0.35),
     [("EVENT PLAN  ·  30 SESSIONS", GOLD, True, FONT_EN)], size=13)
text(s, ML, Inches(2.0), Inches(12), Inches(1.0),
     "创智汇 · 30场活动具体方案", size=34, color=INK, bold=True)
text(s, ML, Inches(3.1), Inches(11), Inches(0.4),
     "总—分—总　｜　主题适配 · 分场策划 · 招商闭环", size=16, color=MUT)
rect(s, ML, Inches(3.7), Inches(2.4), Pt(3), fill=GOLD)
card(s, ML, Inches(4.1), Inches(3.7), Inches(2.2), "园区主题", accent=GOLD,
     items=["AI + 数字内容无界共创港", "3F AI/OPC　5F IP展贸", "同浦汇产业招服执行"])
card(s, Inches(4.7), Inches(4.1), Inches(3.7), Inches(2.2), "交付结构", accent=ACC,
     items=["总：主题与价值", "分：30场逐场策划", "总：协同逼单与闭环"])
card(s, Inches(8.65), Inches(4.1), Inches(3.85), Inches(2.2), "年度基线", accent=GREEN,
     items=["30场 / 打包执行费30万", "每场原则上≥30人", "年触达900+精准客"])

# 02 目录
idx = nxt()
s = slide()
header(s, "00", "AGENDA", "汇报目录 · 总—分—总")
footer(s, idx)
items = [
    ("01 总", "主题定位与价值框架", "为何30场、适配什么、创造什么价值"),
    ("02 总", "六大主题与30场总表", "名称、排期、楼层、类型一览"),
    ("03 分", "六大主题策划案", "行程模板 · 嘉宾 · 落位 · 价值链接"),
    ("04 分", "30场逐场执行卡", "一场一页：行程/嘉宾/排期/转化"),
    ("05 总", "协同支持与招商闭环", "人员配置 · 逼单 · 转化 · 闭环流程"),
]
for i, (n, t, d) in enumerate(items):
    y = Inches(1.45) + Inches(i * 0.95)
    rect(s, ML, y, CW, Inches(0.85),
         grad=[(0, "2A1E55", None), (100, "1E1542", None)], line=LINE, lw=1, radius=True)
    text(s, Emu(ML + Inches(0.3)), Emu(y + Inches(0.15)), Inches(1.4), Inches(0.5),
         [(n, GOLD, True)], size=18)
    text(s, Emu(ML + Inches(1.9)), Emu(y + Inches(0.12)), Inches(8), Inches(0.35),
         t, size=16, color=INK, bold=True)
    text(s, Emu(ML + Inches(1.9)), Emu(y + Inches(0.45)), Inches(9), Inches(0.3),
         d, size=12, color=MUT)

# 03 总·主题适配
idx = nxt()
s = slide()
header(s, "01", "THEME FIT", "总 · 活动如何适配园区主题")
footer(s, idx)
card(s, ML, Inches(1.45), Inches(5.7), Inches(5.2), "园区主题锚点", accent=GOLD, items=[
    "主定位：AI + 数字内容无界共创港",
    "3F：人工智能 / OPC / 孵化办公",
    "5F：IP内容 / 潮玩展贸 / 培训沙龙",
    "政策：YOUNG立方 + 火山引擎 + 服务中心",
    "目标：活动即招商、内容即去化",
])
card(s, Inches(6.85), Inches(1.45), Inches(5.6), Inches(5.2), "六大主题 → 楼层映射", accent=ACC, items=[
    "A AI培训/黑客松 → 主落 3F（10场）",
    "B 政策沙龙/路演 → 3F/5F沙龙（6场）",
    "C 行业对接/撮合 → 主落 5F集群（6场）",
    "D IP/潮玩活动 → 主落 5F展区（4场）",
    "E 漫展/大型 → 5F+临建（2场）",
    "F 领事/出海 → 会客厅/5F（2场）",
])

# 04 总·双价值
idx = nxt()
s = slide()
header(s, "01", "DUAL VALUE", "总 · 给园区的价值 × 给招商的价值")
footer(s, idx)
card(s, ML, Inches(1.45), Inches(5.7), Inches(5.2), "给园区 / 载体的价值", accent=GOLD, items=[
    "完成载体活跃度与活动KPI（满年补贴基础）",
    "提升人气、媒体声量与政府可见度",
    "验证3F/5F主题定位（AI+IP可感知）",
    "沉淀社群、会员与复访数据资产",
    "部分场次创造门票/赞助经营性收入",
    "服务中心/联合会专业形象外溢",
])
card(s, Inches(6.85), Inches(1.45), Inches(5.6), Inches(5.2), "给园区招商的价值", accent=GREEN, items=[
    "每场≥30人精准客，年触达900+",
    "活动中嵌看场/报价/意向表，缩短决策",
    "政策场「带补贴入驻」降低抗性",
    "产业集群场直接去化5F展位",
    "大型场集中留资，会后7日逼单窗口",
    "形成「活动→线索→看场→签约」漏斗",
])

# 05 六大主题总览
idx = nxt()
s = slide()
header(s, "01", "SIX PILLARS", "总 · 六大主题一览")
footer(s, idx)
for i, (k, (name, freq, axis, col, path)) in enumerate(CAT_META.items()):
    x = ML + Inches((i % 3) * 4.0)
    y = Inches(1.45) + Inches((i // 3) * 2.55)
    card(s, x, y, Inches(3.85), Inches(2.4), f"{k}. {name}　{freq}", accent=col, items=[
        f"适配：{axis}",
        f"路径：{path}",
    ], bsize=12)

# 06-07 30场总表
rows1 = [["序号", "主题", "活动名称", "建议档期", "场地", "招商作用"]]
rows2 = [["序号", "主题", "活动名称", "建议档期", "场地", "招商作用"]]
for e in EVENTS:
    cat_name = CAT_META[e["cat"]][0].split(" ")[0]
    row = [str(e["id"]), cat_name, e["name"], e["month"], e["floor"], e["lease"][:18]]
    (rows1 if e["id"] <= 15 else rows2).append(row)
add_table_slide("30场活动总表（1–15）", rows1, [0.07, 0.1, 0.32, 0.14, 0.18, 0.19],
                sizes=[10, 10, 11, 10, 10, 10], rh=0.32)
add_table_slide("30场活动总表（16–30）", rows2, [0.07, 0.1, 0.32, 0.14, 0.18, 0.19],
                sizes=[10, 10, 11, 10, 10, 10], rh=0.32)

# 08 分·主题策划案（6个）
for k, (name, freq, axis, col, path) in CAT_META.items():
    evs = [e for e in EVENTS if e["cat"] == k]
    idx = nxt()
    s = slide()
    header(s, "03", f"THEME {k}", f"分 · 主题策划案｜{name}")
    footer(s, idx)
    card(s, ML, Inches(1.4), Inches(4.0), Inches(5.25), "主题定位与适配", accent=col, items=[
        f"场次：{freq}",
        f"楼层适配：{axis}",
        f"转化路径：{path}",
        f"代表场次：{evs[0]['name']}",
        "统一动作：签到→内容→看场→意向表→回访",
    ], bsize=12)
    # 行程模板
    sample = evs[0]
    card(s, Inches(5.2), Inches(1.4), Inches(7.2), Inches(5.25), "标准行程模板（可套用本主题各场）", accent=GOLD, items=[
        f"建议时长：{sample['dur']}",
        f"落位形式：{sample['format']}　｜　场地：{sample['floor']}",
        f"嘉宾配置：{sample['guests']}",
        *[f"行程：{a}" for a in sample["agenda"][:5]],
        f"现场人力：{sample['staff']}",
    ], bsize=12)

# 09 分·30场逐场（每场1页）
for e in EVENTS:
    idx = nxt()
    s = slide()
    cat_name = CAT_META[e["cat"]][0]
    header(s, "04", f"E{e['id']:02d}", f"{e['id']:02d}　{e['name']}")
    footer(s, idx)
    # meta strip
    rect(s, ML, Inches(1.35), CW, Inches(0.7),
         grad=[(0, "2A1E55", None), (100, "1E1542", None)], line=GOLD, lw=1, radius=True)
    text(s, Emu(ML + Inches(0.25)), Inches(1.42), Inches(11.3), Inches(0.55),
         [(f"{cat_name}　｜　{e['month']}　｜　{e['floor']}　｜　{e['dur']}　｜　{e['format']}　｜　{e['staff']}", GOLD, True)],
         size=13, anchor=MSO_ANCHOR.MIDDLE)

    card(s, ML, Inches(2.2), Inches(4.2), Inches(4.45), "具体行程", accent=GOLD,
         items=e["agenda"], bsize=13)
    card(s, Inches(4.55), Inches(2.2), Inches(3.9), Inches(2.1), "嘉宾与落位", accent=ACC, items=[
        f"嘉宾：{e['guests']}",
        f"形式：{e['format']}",
        f"场地：{e['floor']}",
        "落位：签到台→主区→看场动线→洽谈角",
    ], bsize=12)
    card(s, Inches(8.6), Inches(2.2), Inches(3.95), Inches(2.1), "排期与执行", accent=ACC2, items=[
        f"档期：{e['month']}",
        "T-14：嘉宾锁定+物料",
        "T-7：邀约名单确认",
        "T-1：场地彩排",
        "T+1~7：回访逼单",
    ], bsize=12)
    card(s, Inches(4.55), Inches(4.45), Inches(3.9), Inches(2.2), "园区价值", accent=GREEN, items=[
        e["park"],
        "活跃度/声量/主题验证",
    ], bsize=12)
    card(s, Inches(8.6), Inches(4.45), Inches(3.95), Inches(2.2), "招商价值链接", accent=ROSE, items=[
        e["lease"],
        "看场→意向表→报价→签约",
    ], bsize=12)

# 10 总·协同人员
idx = nxt()
s = slide()
header(s, "05", "STAFFING", "总 · 对方如何协同支持（人员）")
footer(s, idx)
rows = [
    ["场次类型", "我方（同浦汇）", "园区配合", "服务中心/其他", "合计参考"],
    ["常规培训/沙龙", "2–3人", "1–2人（场地/引导）", "0–2人", "4–7人"],
    ["政策诊断/路演", "2–3人", "1–2人", "服务中心2–4人", "5–9人"],
    ["产业集群对接", "3–4人", "2人（展区接待）", "商协会1–2人", "6–8人"],
    ["IP发布/市集", "4–5人", "3–4人", "外包票务视情况", "7–12人"],
    ["漫展/嘉年华", "6–8人", "5–6人+物业安保", "外包执行", "15人+"],
    ["领事/出海", "4–5人", "2人", "外事/翻译协同", "8–10人"],
]
yy = Inches(1.4)
for ri, row in enumerate(rows):
    xx = ML
    h = Inches(0.55)
    cols = [0.18, 0.2, 0.22, 0.22, 0.18]
    for ci, val in enumerate(row):
        cw = Emu(int(CW * cols[ci]))
        bg = RGBColor(0x2A, 0x1E, 0x55) if ri == 0 else (
            RGBColor(0x24, 0x1A, 0x4A) if ri % 2 else RGBColor(0x1C, 0x14, 0x3A))
        rect(s, xx, yy, cw, h, fill=bg, line=LINE, lw=0.5)
        text(s, Emu(xx + Inches(0.08)), Emu(yy + Inches(0.1)), Emu(cw - Inches(0.12)), Emu(h - Inches(0.15)),
             str(val), size=12, color=GOLD if ri == 0 else INK, bold=(ri == 0), anchor=MSO_ANCHOR.MIDDLE)
        xx = Emu(xx + cw)
    yy = Emu(yy + h)
text(s, ML, Inches(5.6), CW, Inches(1.1),
     "园区固定协同清单：①场地与物业报备　②一层导视/停车　③领导是否出席　④意向客户名单共享　⑤会后会议室留给逼单。",
     size=14, color=MUT)

# 11 逼单
idx = nxt()
s = slide()
header(s, "05", "CLOSING", "总 · 怎么样去逼单")
footer(s, idx)
card(s, ML, Inches(1.4), Inches(3.9), Inches(5.25), "场中逼单设计", accent=GOLD, items=[
    "每场必设「看场时段」不少于20分钟",
    "出口处置意向表/政策礼包二维码",
    "限时席位：当场留资享免租/折扣码",
    "嘉宾/导师荐客即时对接招商经理",
    "大型场设「招商绿色通道」专台",
], bsize=13)
card(s, Inches(5.1), Inches(1.4), Inches(3.9), Inches(5.25), "会后7日逼单窗口", accent=ACC, items=[
    "T+1：短信/企微致谢+资料包",
    "T+2：电话回访，约二次看场",
    "T+3–5：政策诊断/报价单发出",
    "T+7：决策人面谈或线上签约",
    "未成交进培育池，导入下场活动",
], bsize=13)
card(s, Inches(9.2), Inches(1.4), Inches(3.35), Inches(5.25), "逼单话术锚点", accent=GREEN, items=[
    "政策窗口期",
    "展位/单元余量",
    "同场竞品已留资",
    "活动专属礼包到期",
    "领导见证签约位",
], bsize=13)

# 12 转化
idx = nxt()
s = slide()
header(s, "05", "CONVERSION", "总 · 怎么样促进转化")
footer(s, idx)
steps = [
    ("线索", "签到·企微·意向表", "场场留资率≥60%"),
    ("培育", "社群·内容·下场邀约", "7日触达≥3次"),
    ("看场", "标准动线+样板间", "看场到访率≥25%"),
    ("报价", "租金+政策礼包组合", "48小时内出单"),
    ("成交", "意向金/合同", "月度转化看板"),
]
for i, (a, b, c) in enumerate(steps):
    x = ML + Inches(i * 2.4)
    rect(s, x, Inches(1.6), Inches(2.25), Inches(2.4),
         grad=[(0, "2A1E55", None), (100, "1E1542", None)], line=LINE, lw=1, radius=True)
    text(s, Emu(x + Inches(0.15)), Inches(1.8), Inches(2), Inches(0.4),
         [(f"0{i+1}", GOLD, True, FONT_EN)], size=20)
    text(s, Emu(x + Inches(0.15)), Inches(2.4), Inches(2), Inches(0.4), a, size=16, color=INK, bold=True)
    text(s, Emu(x + Inches(0.15)), Inches(2.9), Inches(2), Inches(0.8), b, size=12, color=MUT)
    text(s, Emu(x + Inches(0.15)), Inches(3.6), Inches(2), Inches(0.3), c, size=11, color=GREEN, bold=True)

card(s, ML, Inches(4.3), CW, Inches(2.35), "转化加速器（与主题强绑定）", accent=GOLD, items=[
    "AI场：工位/小单元「体验日」当天可试坐 → 当周签约有礼",
    "政策场：诊断结论页直接附「入驻可申报清单」→ 带政策决策",
    "集群场：展位扣点测算表现场出 → 降低纯租金抗性",
    "大型场：门票票根/展商证兑换看场礼 → 流量变线索",
], bsize=13)

# 13 闭环
idx = nxt()
s = slide()
header(s, "05", "CLOSED LOOP", "总 · 后期招商闭环流程")
footer(s, idx)
loop = [
    ("1.活动前", "名单筛选·邀约·政策礼包预热·招商目标场次化"),
    ("2.活动中", "签到建档·内容种草·看场·意向表·限时礼"),
    ("3.活动后", "7日逼单·报价·二次看场·未成交进池"),
    ("4.成交", "合同·付款·入驻排期·服务中心接棒申报"),
    ("5.复盘", "场次ROI·转化漏斗·内容沉淀·下场优化"),
    ("6.反哺", "租户故事→下场案例；载体KPI→补贴；社群→下场观众"),
]
for i, (t, d) in enumerate(loop):
    y = Inches(1.4) + Inches(i * 0.85)
    rect(s, ML, y, CW, Inches(0.75),
         grad=[(0, "2A1E55", None), (100, "1E1542", None)], line=LINE, lw=0.8, radius=True)
    text(s, Emu(ML + Inches(0.3)), Emu(y + Inches(0.2)), Inches(2.2), Inches(0.4),
         [(t, GOLD, True)], size=15)
    text(s, Emu(ML + Inches(2.6)), Emu(y + Inches(0.22)), Inches(9), Inches(0.4),
         d, size=14, color=INK)

# 14 结语
idx = nxt()
s = slide()
header(s, "END", "SUMMARY", "结论 · 给对方的协同清单")
footer(s, idx)
card(s, ML, Inches(1.4), Inches(5.7), Inches(5.25), "我们交付", accent=GOLD, items=[
    "30场具名活动+年度排期",
    "六大主题与园区AI+IP主轴对齐",
    "逐场行程/嘉宾/落位/价值链接",
    "执行团队与物料（活动执行费口径）",
    "会后7日逼单与转化漏斗运营",
], bsize=14)
card(s, Inches(6.85), Inches(1.4), Inches(5.6), Inches(5.25), "请园区协同", accent=GREEN, items=[
    "按场次类型配备1–6名现场支持",
    "共享意向客户与在谈名单",
    "开放样板间/展区看场动线",
    "确认领导出席与媒体口径",
    "会后会议室与逼单权限支持",
    "月度联席复盘转化数据",
], bsize=14)

# fix footers
total = len(_FOOT)
for run in _FOOT:
    run.text = " / %02d" % total

# save ppt
os.makedirs(ART, exist_ok=True)
out_cn = os.path.join(HERE, "创智汇30场活动具体方案.pptx")
out_en = os.path.join(HERE, "chuangzhihui-30-events-plan.pptx")
prs.save(out_cn)
prs.save(out_en)
import shutil
shutil.copy2(out_cn, os.path.join(ART, "创智汇30场活动具体方案.pptx"))
shutil.copy2(out_en, os.path.join(ART, "chuangzhihui-30-events-plan.pptx"))
print("PPT slides:", len(prs.slides), "->", out_cn)

# Excel sync
wb = Workbook()
ws = wb.active
ws.title = "30场总表"
headers = ["序号", "主题代码", "主题类型", "活动名称", "建议档期", "场地", "时长", "形式",
           "嘉宾", "行程要点", "园区价值", "招商价值", "现场人力", "逼单要点"]
for j, h in enumerate(headers, 1):
    cell = ws.cell(1, j, h)
    cell.font = Font(name="微软雅黑", bold=True, color="FFFFFF")
    cell.fill = PatternFill("solid", fgColor="4A2C7A")
    cell.alignment = Alignment(wrap_text=True, vertical="center")
for i, e in enumerate(EVENTS, 2):
    vals = [
        e["id"], e["cat"], CAT_META[e["cat"]][0], e["name"], e["month"], e["floor"],
        e["dur"], e["format"], e["guests"], " / ".join(e["agenda"]),
        e["park"], e["lease"], e["staff"], "场中看场+意向表；T+1~7回访逼单",
    ]
    for j, v in enumerate(vals, 1):
        cell = ws.cell(i, j, v)
        cell.font = Font(name="微软雅黑", size=9)
        cell.alignment = Alignment(wrap_text=True, vertical="center")
        if i % 2 == 0:
            cell.fill = PatternFill("solid", fgColor="F3E9FF")
ws.row_dimensions[1].height = 28
for i in range(2, 32):
    ws.row_dimensions[i].height = 48
for i, w in enumerate([6, 8, 16, 28, 12, 16, 14, 12, 22, 40, 28, 28, 16, 24], 1):
    ws.column_dimensions[get_column_letter(i)].width = w

ws2 = wb.create_sheet("协同与闭环")
for i, row in enumerate([
    ["环节", "动作", "园区协同", "成功标准"],
    ["活动前", "名单筛选、邀约、礼包预热", "共享在谈客户、确认场地物业", "到会名单锁定"],
    ["活动中", "签到建档、看场、意向表、限时礼", "1–6人现场支持（按场次）", "留资率≥60%"],
    ["会后7日", "回访、报价、二次看场", "会议室/决策人对接", "看场率≥25%"],
    ["成交", "合同、入驻、申报接棒", "合同与物业交割", "月度签约看板"],
    ["复盘反哺", "ROI与漏斗优化", "月度联席会", "下场转化提升"],
], 1):
    for j, v in enumerate(row, 1):
        cell = ws2.cell(i, j, v)
        cell.font = Font(name="微软雅黑", bold=(i == 1), color=("FFFFFF" if i == 1 else "2C2C2C"))
        if i == 1:
            cell.fill = PatternFill("solid", fgColor="4A2C7A")
        cell.alignment = Alignment(wrap_text=True, vertical="center")
for i, w in enumerate([12, 36, 32, 18], 1):
    ws2.column_dimensions[get_column_letter(i)].width = w

x_cn = os.path.join(HERE, "创智汇30场活动方案排期表.xlsx")
x_en = os.path.join(HERE, "chuangzhihui-30-events-schedule.xlsx")
wb.save(x_cn)
wb.save(x_en)
shutil.copy2(x_cn, os.path.join(ART, "创智汇30场活动方案排期表.xlsx"))
shutil.copy2(x_en, os.path.join(ART, "chuangzhihui-30-events-schedule.xlsx"))
print("Excel ->", x_cn)
