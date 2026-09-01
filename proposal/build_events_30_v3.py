# -*- coding: utf-8 -*-
"""创智汇 · 30场活动方案 V3（表格化应答版）
回应甲方三点：
1) 嘉宾/来访者与孵化器关联度（明确KPI，非口号100%）
2) 执行分工（我方做什么 / 园区支持什么）
3) 各档往期案例（可见企业 + 议程 + 现场照片）
并补全：活动概况/价值/交付/收费；租金3.3、物业13.8、待租8间、佣金2个月。
用法: python3 build_events_30_v3.py
"""
import os, shutil
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

HERE = os.path.dirname(os.path.abspath(__file__))
ART = "/opt/cursor/artifacts"
PHOTO = os.path.join(HERE, "cases", "photos")
FONT, FONT_EN = "Microsoft YaHei", "Arial"
BG_A, BG_B = "1A1140", "3A206E"
INK = RGBColor(0xF5, 0xF1, 0xFF)
MUT = RGBColor(0xC6, 0xBA, 0xEA)
SOFT = RGBColor(0x94, 0x86, 0xC4)
ACC = RGBColor(0x8B, 0x7B, 0xFF)
ACC2 = RGBColor(0xC7, 0x7D, 0xFF)
GOLD = RGBColor(0xE9, 0xC2, 0x7C)
GREEN = RGBColor(0x56, 0xD6, 0xB6)
ROSE = RGBColor(0xF5, 0x8B, 0xB8)
LINE = RGBColor(0x4A, 0x3C, 0x7C)
SW, SH = Inches(13.333), Inches(7.5)
ML, CW = Inches(0.55), Inches(12.23)
prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
BLANK = prs.slide_layouts[6]; _FOOT = []; IDX = 0

# ---------- 商业锚点（与本次沟通一致） ----------
COMMERCIAL = {
    "rent": "办公租金 3.3 元/㎡/天",
    "property": "物业费 13.8 元/㎡/月",
    "rooms": "待租办公室 8 间（优先活动转化去化）",
    "commission": "招商佣金：成交首年有效净租金 × 2 个月（首年不重复）",
    "activity_fee": "活动打包执行费：30 场 / 年 打包 30 万元（含策划执行、基础物料、主持协调）",
    "media_fee": "媒体流量费：约 10 万元/年（另计，可冲抵活动获客）",
}

# ---------- 关联度口径（回答“是不是100%”） ----------
ASSOC_RULE = {
    "invite_target": "邀约名单与孵化器画像关联度 ≥90%",
    "onsite_target": "到场核验关联度 ≥80%",
    "dm_note": "企业负责人（法人/CEO/创始人/业务VP/采购决策人）占比按档位设定 35%–60%",
    "not_100": "不是简单承诺“到场100%都是入孵企业”：允许约10%–15%生态角色（政府/高校/投资人/媒体），用于背书与转化；但主体客群必须是 AI/OPC/数字内容等可入驻画像企业。",
    "if_require_100": "若甲方要求“除官方嘉宾外，来访企业名单100%画像匹配”，可按场次书面锁定邀约名单审核。",
}

CAT_META = {
    "A": ("智能体/训练营", "8场", "Agent·多模态·算力·具身", GOLD, "40–80人", "35–45%", "≥90%"),
    "B": ("政策/治理沙龙", "5场", "政策兑现·高企·三券", ACC, "40–80人", "50–60%", "≥90%"),
    "C": ("产业/科学对接", "5场", "成果转化·集群对接", ACC2, "50–100人", "55–65%", "≥85%"),
    "D": ("内容IP活动", "4场", "首发·市集·联展", GREEN, "80–300人", "30–40%", "≥75%"),
    "E": ("WAIC联动大场", "4场", "开放日+双夜+嘉年华", ROSE, "100–400人", "40–55%", "≥80%"),
    "F": ("出海/AGI/收官", "4场", "买家团·领事·年终", GOLD, "40–120人", "50–60%", "≥90%"),
}

# 执行分工模板
DIV_US = "选题策划；嘉宾主邀；物料/主持；签到建档；场控；会后7日回访逼单；复盘报告与意向分级"
DIV_PARK = "场地物业报备；水电空调；导视停车；样板间开放；在谈客户名单共享；现场支持1–N人；合同/交割对接"
DIV_UNION = "会员协同邀约；政策宣讲（服中心）；领导出席协调；申报服务接棒"

# 交付物模板
DELIVER = "①活动方案+议程 ②邀约名单与关联度核验表 ③签到表/企微建档 ④现场执行 ⑤意向客户分级表 ⑥7日回访纪要 ⑦单场复盘PPT/一页纸"

FEE_PACK = "计入年度30万打包（本场约分摊见收费表）；门票/赞助另计归属约定；转化成交另计佣金2个月"

RHYTHM = "T-21选题定档 → T-14嘉宾物料 → T-7邀约确认 → T-1彩排 → D日执行 → T+1名单 → T+7回访逼单 → T+14复盘"

def ev(id_, cat, waic, name, month, floor, dur, fmt, people, dm, assoc, guests, agenda,
       value, deliver=DELIVER, us=DIV_US, park=DIV_PARK, fee=FEE_PACK, case_tier=None,
       budget_share=1.0):
    return dict(
        id=id_, cat=cat, waic=waic, name=name, month=month, floor=floor, dur=dur, format=fmt,
        people=people, dm=dm, assoc=assoc, guests=guests, agenda=agenda, value=value,
        deliver=deliver, us=us, park=park, union=DIV_UNION, fee=fee,
        case_tier=case_tier or cat, rhythm=RHYTHM, budget_share=budget_share,
    )

EVENTS = [
    ev(1,"A","智能伙伴","智能伙伴·创智汇开年主题沙龙","1月下旬","5F沙龙","2.5小时","主题沙龙","50–70人","45%","≥90%",
       "联合会领导、AI/OPC企业负责人、服中心",
       "18:30签到参观→19:00年度主题解读→19:30园区定位发布→20:10圆桌→20:40看场意向",
       "开年定调；集中获客；8间待租首波曝光", budget_share=0.8),
    ev(2,"A","Agent","Agent智能体搭建一日营","2月中旬","3F OPC","1天","实训营","40–60人","40%","≥90%",
       "Agent工程师、开源Maintainer、初创CEO",
       "09:30签到→10:00工作流拆解→13:30低代码实操→16:00Demo→16:40看场逼单",
       "Agent团队入孵转化；OPC工位去化", budget_share=1.0),
    ev(3,"A","多模态","多模态智能体工作坊","3月上旬","3F培训+直播间","1天","工作坊","40–60人","40%","≥90%",
       "多模态研究者、应用厂商产品负责人",
       "案例分享→视觉语言实操→场景共创→作品评审→政策礼包+看场",
       "多模态应用团队入驻；内容工具链粘性", budget_share=1.0),
    ev(4,"A","算力Infra","火山引擎×算力Infra实务营","3月下旬","3F+云创基地","1天","厂商联训","50–80人","45%","≥90%",
       "火山引擎讲师、云创基地、高算力企业CTO",
       "算力政策包→代金券实操→集群案例→一对一诊断→入驻捆绑云资源",
       "高算力消耗企业定向招商；云创协同", budget_share=1.2),
    ev(5,"A","具身智能","具身智能空间交互体验日","4月上旬","5F展区+3F","半天","体验日","40–60人","40%","≥90%",
       "机器人/具身团队、高校实验室负责人",
       "具身演示→空间交互讲解→场景清单→5F踩点→落位洽谈",
       "具身/机器人团队看场转化", budget_share=0.9),
    ev(6,"A","AI内容","AIGC微短剧制片特训","5月上旬","3F直播间+5F","1天","特训营","40–60人","35%","≥90%",
       "微短剧导演、厂牌负责人、平台方",
       "YOUNG立方政策→制片流程→脚本工作坊→路演→5F看场",
       "内容厂牌/工作室入驻；5F场景使用", budget_share=1.0),
    ev(7,"A","Builders","OPC超级个体黑客松（春）","5月中旬","3F整层","1.5天","黑客松","60–80人","40%","≥90%",
       "评委×3、投资人×2、云厂商、Builders",
       "开题组队→封闭开发→路演决赛→入驻礼包→一对一看场",
       "获奖团队优先谈单元；青年创业种子库", budget_share=1.5),
    ev(8,"A","Agent","AI营销Agent实战营","10月中旬","3F培训","半天","实战营","40–50人","40%","≥90%",
       "投放操盘手、Agent产品经理、营销公司负责人",
       "投放策略→Agent自动化→素材AIGC→复盘→服务包转化",
       "营销型公司入驻；服务包续费", budget_share=0.8),
    ev(9,"B","智能伙伴","YOUNG立方×智能伙伴政策沙龙","1月中旬","5F沙龙","2.5小时","政策沙龙","50–80人","55%","≥90%",
       "区政策宣讲、服中心、内容/AI企业负责人",
       "政策要点→园区礼包→适用画像→诊断预约→看场",
       "政策前置获客；内容/AI导入", budget_share=0.8),
    ev(10,"B","治理","AI治理与可信智能体沙龙","2月下旬","3F","2小时","闭门/沙龙","40–60人","55%","≥90%",
       "治理学者、合规顾问、企业法务/CEO",
       "治理框架→合规清单→案例→问答→入驻服务对接",
       "信任背书；合规意识型企业入驻", budget_share=0.8),
    ev(11,"B","产业","高企认定冲刺（AI企业专场）","3月中旬","3F","半天","辅导会","40–60人","60%","≥90%",
       "高企辅导、财税顾问、待认定AI企业负责人",
       "条件拆解→材料要点→初筛→辅导套餐→入驻激励",
       "外区待认定企业带政策入驻", budget_share=0.9),
    ev(12,"B","产业","专精特新·AI应用培育路演","6月上旬","5F","半天","路演","50–80人","55%","≥90%",
       "评审顾问、银行、基金、成长型AI企业",
       "梯度政策→企业路演→点评→金融对接→看场",
       "优质企业筛选补位", budget_share=1.0),
    ev(13,"B","算力","创新券·算力券·模型券实务沙龙","11月上旬","3F","2小时","实务沙龙","40–70人","50%","≥90%",
       "创新券平台、云厂商、用券企业负责人",
       "三券规则→核销实操→案例→开户引导→入驻转化",
       "用券企业向园区聚集；服中心KPI", budget_share=0.8),
    ev(14,"C","AI for Science","高校成果转化·AI for Science日","4月中旬","3F+沙龙","半天","对接日","50–80人","55%","≥85%",
       "复旦/同济技术转移、教授团队、成果公司",
       "成果路演→科学智能体场景→园区承接→看孵化单元→转化落位",
       "成果公司/实验室落户", budget_share=1.0),
    ev(15,"C","产业","汕头玩具×AIGC供应链对接会","4月下旬","5F玩具区","半天","对接会","60–100人","60%","≥85%",
       "汕头商协会、品牌采购、玩具企业负责人",
       "集群介绍→AIGC设计提效→供需对接→展位参观→报价",
       "5F玩具展位去化", budget_share=1.0),
    ev(16,"C","产业","扬州毛绒×数字人联名对接","6月中旬","5F毛绒区","半天","对接会","50–80人","60%","≥85%",
       "扬州集群、数字人厂商、品牌方",
       "联名模式→租金扣点测算→看场→MOU→跟进",
       "700㎡量级客户线索", budget_share=1.0),
    ev(17,"C","产业","东莞潮玩品牌入沪推介","8月上旬","5F潮玩区","半天","推介会","50–80人","55%","≥85%",
       "潮玩品牌、渠道商负责人",
       "上海渠道→5F落位→联名活动→看场→定金意向",
       "品牌展位/快闪去化", budget_share=1.0),
    ev(18,"C","内容","IP授权×AI衍生交易撮合会","9月中旬","5F展示中心","1天","撮合会","60–100人","55%","≥85%",
       "IP方、被授权商、律师、渠道",
       "IP路演→AI衍生专题桌→一对一撮合→成交看板→落户激励",
       "IP/衍生品团队入驻；贸易成交", budget_share=1.2),
    ev(19,"D","内容","AI创作者内容首发①","5月下旬","5F主展","3小时","发布会","80–150人","35%","≥75%",
       "创作者、媒体、渠道、品牌负责人",
       "揭幕→发布→签售快闪→媒体专访→招商通道",
       "内容品牌曝光；关联品牌问询", budget_share=1.2),
    ev(20,"D","内容","潮玩×AIGC主题市集","6月下旬","5F+公区","2天","市集","200–400人","30%","≥75%",
       "摊主×40、达人、潮玩品牌",
       "布摊→开市→AIGC打卡→夜间场→优质摊主转正谈",
       "人气经营；摊主升级固定展位/办公", budget_share=1.5),
    ev(21,"D","内容","沉浸式AI+IP联展","8月中旬","5F展示中心","7–10天","联展","日均80–200人","35%","≥75%",
       "联合IP×5、文旅渠道、品牌",
       "布展→开幕→预约观展→教育场→闭幕招商酒会",
       "长期人流；闭幕酒会集中转化", budget_share=1.8),
    ev(22,"D","内容","AI创作者内容首发②","10月下旬","5F","3小时","发布会","80–150人","35%","≥75%",
       "IP方、渠道、媒体",
       "发布→双十一预售→渠道对接→看场→年框意向",
       "旺季补位；年框意向", budget_share=1.2),
    ev(23,"E","WAIC主会","WAIC UP! 创智汇「智能伙伴」开放日","7月17日","3F+5F整区","1天","开放日","150–300人","45%","≥80%",
       "WAIC回流嘉宾、媒体、AI/内容企业、高校",
       "10:00致辞→10:30主题速递→11:30开放参观→14:00双展台→15:30政策礼包→16:30意向洽谈",
       "大会流量导入；8间待租集中看场", budget_share=2.0),
    ev(24,"E","Builders夜","WAIC UP! AI Builders Night","7月17日晚","5F主厅","18:30–21:30","夜场社交","80–150人","45%","≥85%",
       "创业者、开发者、投资人、开源Maintainer",
       "签到名牌→致辞→闪电演讲×4→开放麦→AMA→投资人蹲项目+看场",
       "建造者当场看3F单元", budget_share=1.8),
    ev(25,"E","AGI夜","WAIC UP! 通往AGI之夜","7月18日晚","5F圆桌区","19:00–21:30","学术夜场","60–100人","50%","≥85%",
       "复旦学者、大模型负责人、Agent创始人、治理学者",
       "复旦致辞→无PPT圆桌→三桌围炉→辩题投票→深度社交+精品看场",
       "高端研究型/涉外企业入驻", budget_share=1.8),
    ev(26,"E","青年菁英","创智汇秋季潮玩×AI嘉年华","10月上旬","5F+公区","2天","嘉年华","300–600人","35%","≥75%",
       "潮玩品牌、达人、赞助商、AI创作者",
       "开幕秀→市集展洽→AI赛事→夜间场→签约仪式",
       "集中签约展位/办公", budget_share=2.2),
    ev(27,"F","出海","AI出海推介与国际买家团","9月下旬","5F+培训","1天","买家团","50–80人","55%","≥90%",
       "海外买家、跨境平台、出海企业负责人",
       "出海路径→买家需求→一对一洽谈→订单意向→办公/展位捆绑",
       "出海需求企业入驻；撮合收入", budget_share=1.2),
    ev(28,"F","AGI","通往AGI季度圆桌①","4月下旬","5F沙龙","2.5小时","圆桌","40–60人","55%","≥90%",
       "学者、模型企业、投资人",
       "议题导入→圆桌→围炉→问答→看场",
       "研究型/模型团队入驻线索", budget_share=0.8),
    ev(29,"F","领事","国别会客厅·AI合作领事专题","11月中旬","会客厅/5F","2.5小时","领事沙龙","40–70人","55%","≥90%",
       "领事官员、涉外AI企业负责人、翻译",
       "外事接待→国别AI机会→B2B闪见→国际服务→高端看场",
       "涉外/出海企业入驻背书", budget_share=1.2),
    ev(30,"F","智能伙伴","创智汇AI年度Demo Day·收官","12月上旬","5F主展","1天","路演日","80–150人","50%","≥90%",
       "投资人、链主、媒体、区级嘉宾、十强项目",
       "开幕→十强路演→颁奖→明年WAIC预热→红酒洽谈+集中签约",
       "集中签约；来年预热；媒体背书", budget_share=1.5),
]

# 往期案例（公开活动，含可见企业/机构 + 议程 + 照片）
CASES = [
    {
        "tier": "B/F·政策治理/出海",
        "name": "「杨数浦」数字经济出海合规沙龙（第三期）",
        "when": "2025-02-18",
        "where": "上海人民城市实践展示馆（杨浦滨江）",
        "orgs": "杨浦区委网信办、区工商联；赛博网络安全产业创新研究院；上海社科院互联网研究中心",
        "companies": "惠志斌（社科院）及出海合规顾问团队；数字经济企业管理层/法务/技术代表",
        "agenda": [
            "主题：数字新征程，出海创未来——合规引领发展",
            "专家分享：数据跨境与出海合规策略",
            "圆桌：企业出海方向与合规路径抉择",
            "授证：出海合规荣誉顾问授证仪式",
        ],
        "photos": ["case_65931516.jpg", "case_40a029da.jpg", "case_fb37cc44.jpg"],
        "source": "上海市杨浦区人民政府官网报道",
        "map_to": "对应创智汇：政策沙龙(B)、出海买家团/领事专题(F)",
    },
    {
        "tier": "B·政策/治理",
        "name": "杨「数」浦数字沙龙第七期：AI如何重塑企业DNA",
        "when": "2025-06-30",
        "where": "美团上海综合指挥中心",
        "orgs": "杨浦区委网信办主办；赛博院、上海市人工智能与社会发展研究会承办",
        "companies": "美团；大众点评大模型应用团队；赛博院研究员卢虹羽等",
        "agenda": [
            "参观美团综合指挥中心数字化场景",
            "主题分享：人机共生时代·AI重塑企业DNA",
            "案例：大众点评AI技术应用创新与落地",
            "政策解读：《人工智能生成合成内容标识办法》等",
            "互动答疑",
        ],
        "photos": ["case_c987bd91.jpg", "case_59157ea5.jpg", "case_6e51af9e.jpg", "case_76d39051.jpg"],
        "source": "上海市杨浦区人民政府官网报道",
        "map_to": "对应创智汇：AI治理沙龙、智能伙伴政策沙龙、Agent应用营",
    },
    {
        "tier": "B/C·路演对接",
        "name": "「融见科创·智启未来」人工智能专场路演暨投融资对接会",
        "when": "2025-10（杨浦科创促进会）",
        "where": "杨浦（邮储银行联合主办）",
        "orgs": "杨浦科创促进会、邮储银行上海分行；科辰创投、智能产业创新研究院、益华资本等",
        "companies": "复楚智能、卡房信息、一造科技、中科趋势、万笔千墨 等路演企业",
        "agenda": [
            "开场致辞（促进会会长夏立城）",
            "主办方致辞（邮储银行）",
            "主题分享：AI行业趋势（中邮证券）",
            "“6+6+6+1”模式路演对接（6项目+6投资人+6金融机构+1政府代表）",
            "区投促办代表点评与招商邀请",
            "评委总结与后续跟进",
        ],
        "photos": ["case_db1d5cac.jpg", "case_cb0a3812.jpg", "case_cdfd470a.jpg"],
        "source": "上海市杨浦区人民政府官网 / 青年报报道",
        "map_to": "对应创智汇：专精特新路演、Demo Day、黑客松决赛日",
    },
]


def _grad(shape, stops, ang=90):
    spPr = shape._element.spPr
    for t in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill"):
        e = spPr.find(qn(t))
        if e is not None: spPr.remove(e)
    g = spPr.makeelement(qn("a:gradFill"), {}); lst = g.makeelement(qn("a:gsLst"), {})
    for pos, col, al in stops:
        gs = g.makeelement(qn("a:gs"), {"pos": str(int(pos * 1000))})
        c = g.makeelement(qn("a:srgbClr"), {"val": col})
        if al is not None: c.append(g.makeelement(qn("a:alpha"), {"val": str(int(al * 1000))}))
        gs.append(c); lst.append(gs)
    g.append(lst); g.append(g.makeelement(qn("a:lin"), {"ang": str(int(ang * 60000)), "scaled": "1"}))
    ln = spPr.find(qn("a:ln")); (ln.addprevious(g) if ln is not None else spPr.append(g))


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.shadow.inherit = False; bg.line.fill.background()
    _grad(bg, [(0, BG_A, None), (55, BG_B, None), (100, BG_A, None)], 120)
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s


def rect(s, x, y, w, h, fill=None, line=None, lw=1.0, radius=False, grad=None, gang=90):
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    b.shadow.inherit = False
    if grad is not None: _grad(b, grad, gang)
    elif fill is None: b.fill.background()
    else: b.fill.solid(); b.fill.fore_color.rgb = fill
    if line is None: b.line.fill.background()
    else: b.line.color.rgb = line; b.line.width = Pt(lw)
    return b


def text(s, x, y, w, h, runs, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=1.0, font=FONT):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    if isinstance(runs, str): runs = [(runs, color, bold)]
    p = tf.paragraphs[0]; p.alignment = align; p.line_spacing = space
    for seg in runs:
        t, c, b = seg[0], seg[1], seg[2]; f = seg[3] if len(seg) > 3 else font
        r = p.add_run(); r.text = t; r.font.size = Pt(size); r.font.bold = b; r.font.color.rgb = c; r.font.name = f
    return tb


def bullets(s, x, y, w, h, items, size=12, color=MUT, gap=4, mark=ACC):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.space_after = Pt(gap); p.line_spacing = 1.12
        r0 = p.add_run(); r0.text = "▪  "; r0.font.size = Pt(size); r0.font.color.rgb = mark; r0.font.name = FONT
        r = p.add_run(); r.text = it; r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = FONT
    return tb


def header(s, sec, eyebrow, title_):
    rect(s, ML, Inches(0.38), Inches(0.06), Inches(0.7), fill=GOLD)
    text(s, Emu(ML + Inches(0.18)), Inches(0.36), Inches(10), Inches(0.26), eyebrow, size=11, color=GOLD, bold=True, font=FONT_EN)
    text(s, Emu(ML + Inches(0.18)), Inches(0.62), Inches(11), Inches(0.42), title_, size=20, color=INK, bold=True)
    text(s, Inches(11.5), Inches(0.35), Inches(1.4), Inches(0.65), [(sec, GOLD, True, FONT_EN)], size=26, align=PP_ALIGN.RIGHT, font=FONT_EN)
    rect(s, ML, Inches(7.05), CW, Pt(1), fill=LINE)


def footer(s, idx):
    text(s, ML, Inches(7.1), Inches(9.5), Inches(0.26),
         [("上海创智汇 ", SOFT, False), ("· 30场表格化方案V3 · 同浦汇", SOFT, False)], size=9)
    rs = text(s, Inches(11.0), Inches(7.1), Inches(1.8), Inches(0.26),
              [("%02d" % idx, GOLD, True, FONT_EN), (" / XX", SOFT, False, FONT_EN)], size=10, align=PP_ALIGN.RIGHT)
    _FOOT.append(rs.text_frame.paragraphs[0].runs[1])


def card(s, x, y, w, h, title=None, items=None, body=None, accent=ACC, tsize=13, bsize=11):
    rect(s, x, y, w, h, grad=[(0, "2A1E55", None), (100, "1E1542", None)], gang=120, line=LINE, lw=1, radius=True)
    rect(s, x, y, Inches(0.07), h, fill=accent)
    if title:
        text(s, Emu(x + Inches(0.22)), Emu(y + Inches(0.12)), Emu(w - Inches(0.35)), Inches(0.34), title, size=tsize, color=INK, bold=True)
    if body:
        text(s, Emu(x + Inches(0.22)), Emu(y + Inches(0.48)), Emu(w - Inches(0.4)), Emu(h - Inches(0.58)), body, size=bsize, color=MUT, space=1.12)
    if items:
        bullets(s, Emu(x + Inches(0.22)), Emu(y + Inches(0.46)), Emu(w - Inches(0.4)), Emu(h - Inches(0.55)), items, size=bsize, mark=accent)


def nxt():
    global IDX; IDX += 1; return IDX


def table_block(s, y0, rows, col_w, sizes=None, rh=0.32, max_h=5.4):
    yy = y0
    if sizes is None: sizes = [10] * len(rows[0])
    for ri, row in enumerate(rows):
        xx = ML; h = Inches(0.38 if ri == 0 else rh)
        for ci, val in enumerate(row):
            cw = Emu(int(CW * col_w[ci]))
            bg = RGBColor(0x2A, 0x1E, 0x55) if ri == 0 else (RGBColor(0x24, 0x1A, 0x4A) if ri % 2 else RGBColor(0x1C, 0x14, 0x3A))
            rect(s, xx, yy, cw, h, fill=bg, line=LINE, lw=0.4)
            text(s, Emu(xx + Inches(0.04)), Emu(yy + Inches(0.03)), Emu(cw - Inches(0.06)), Emu(h - Inches(0.05)),
                 str(val), size=sizes[ci], color=GOLD if ri == 0 else INK, bold=(ri == 0), anchor=MSO_ANCHOR.MIDDLE)
            xx = Emu(xx + cw)
        yy = Emu(yy + h)
    return yy


def add_pic(s, path, x, y, w, h):
    if not os.path.exists(path):
        rect(s, x, y, w, h, fill=RGBColor(0x2A, 0x1E, 0x55), line=LINE, radius=True)
        text(s, x, y, w, h, "照片待补", size=12, color=SOFT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return
    try:
        s.shapes.add_picture(path, x, y, width=w, height=h)
    except Exception:
        rect(s, x, y, w, h, fill=RGBColor(0x2A, 0x1E, 0x55), line=LINE)
        text(s, x, y, w, h, "照片加载失败", size=11, color=SOFT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ===================== SLIDES =====================
s = slide()
text(s, ML, Inches(1.2), Inches(12), Inches(0.3), [("TABLE-FIRST  ·  EVENT OPS V3", GOLD, True, FONT_EN)], size=12)
text(s, ML, Inches(1.65), Inches(12), Inches(0.9), "创智汇 · 30场活动怎么做", size=32, color=INK, bold=True)
text(s, ML, Inches(2.6), Inches(12), Inches(0.4), "表格说清楚：概况 · 价值 · 交付 · 分工 · 收费 · 往期案例", size=15, color=GOLD, bold=True)
text(s, ML, Inches(3.1), Inches(12), Inches(0.35), "同步商业条款：办公3.3元/㎡/天 · 物业13.8元/㎡/月 · 待租8间 · 佣金2个月", size=13, color=MUT)
card(s, ML, Inches(3.7), Inches(3.9), Inches(2.9), "答疑① 关联度", accent=GOLD, items=[
    "邀约画像匹配 ≥90%", "到场核验 ≥80%", "负责人占比按档 35–60%", "不是口号式100%",
], bsize=12)
card(s, Inches(4.7), Inches(3.7), Inches(3.9), Inches(2.9), "答疑② 分工", accent=ACC, items=[
    "我方：策划邀约执行回访", "园区：场地物业名单看场", "联合会/服中心：政策协同", "一张RACI表写死",
], bsize=12)
card(s, Inches(8.85), Inches(3.7), Inches(3.9), Inches(2.9), "答疑③ 案例", accent=GREEN, items=[
    "杨数浦出海合规沙龙", "AI重塑企业DNA沙龙", "融见科创AI路演", "含企业名+议程+现场照",
], bsize=12)

# 商业条款
idx = nxt(); s = slide(); header(s, "00", "COMMERCIAL", "同步商业条款（与本次沟通一致）"); footer(s, idx)
rows = [
    ["条款", "约定", "说明"],
    ["办公租金", "3.3 元/㎡/天", "3F办公/孵化单元报价锚点"],
    ["物业费", "13.8 元/㎡/月", "按月计收"],
    ["待租库存", "办公室 8 间", "活动优先导流看场与逼单"],
    ["招商佣金", "成交首年有效净租金 × 2 个月", "首年不重复；活动转化成交同样适用"],
    ["活动收费", "30 场打包 30 万元/年", "含策划执行、基础物料、主持协调"],
    ["媒体（可选）", "约 10 万元/年", "投流获客，另计"],
]
table_block(s, Inches(1.25), rows, [0.22, 0.38, 0.4], sizes=[12, 13, 12], rh=0.55)
text(s, ML, Inches(5.6), CW, Inches(1.0),
     "活动价值最终落到：8间去化 + 租金回收 + 载体KPI/补贴 + 品牌声量；佣金按成交计，不对出租率对赌。", size=13, color=MUT)

# 关联度专项答疑
idx = nxt(); s = slide(); header(s, "01", "INCUBATOR FIT", "答疑①：嘉宾/来访者与孵化器关联度——不是口号100%"); footer(s, idx)
card(s, ML, Inches(1.25), Inches(6.0), Inches(5.5), "明确KPI（可写进合同附件）", accent=GOLD, items=[
    ASSOC_RULE["invite_target"],
    ASSOC_RULE["onsite_target"],
    ASSOC_RULE["dm_note"],
    "核验方式：报名表行业标签 + 名片/企微 + 当场抽检",
    "月度联席会公示关联度达成率",
], bsize=13)
card(s, Inches(7.25), Inches(1.25), Inches(5.5), Inches(5.5), "为什么不是简单说100%", accent=ROSE, items=[
    ASSOC_RULE["not_100"],
    ASSOC_RULE["if_require_100"],
    "大场（市集/嘉年华）流量属性更强，关联度下限≥75%，但看场通道仍只对画像企业开放",
    "每场输出《关联度核验表》作为交付物",
], bsize=12)

# 六档标准
idx = nxt(); s = slide(); header(s, "01", "TIER KPI", "六大档位标准：人数 · 负责人占比 · 孵化器关联度"); footer(s, idx)
rows = [["档位", "场次", "适配主题", "人数", "负责人占比", "孵化器关联度"]]
for k, (name, freq, axis, col, ppl, dm, assoc) in CAT_META.items():
    rows.append([f"{k}.{name}", freq, axis, ppl, dm, assoc])
table_block(s, Inches(1.25), rows, [0.22, 0.1, 0.28, 0.14, 0.13, 0.13], sizes=[11, 11, 11, 11, 11, 11], rh=0.58)

# 分工RACI
idx = nxt(); s = slide(); header(s, "02", "RACI", "答疑②：执行分工——哪些我们做，哪些需要你们支持"); footer(s, idx)
rows = [
    ["环节", "同浦汇（我方）", "园区/科创集团（贵方）", "联合会/服中心"],
    ["选题定档", "主责：年度日历与单场方案", "确认档期与禁办日", "政策窗口协同"],
    ["嘉宾邀约", "主责：企业/导师/投资人", "共享在谈客户名单", "会员协同邀约"],
    ["场地物料", "主责：活动物料/签到/主持", "主责：场地物业水电导视停车", "领导出席协调"],
    ["现场执行", "主责：流程场控建档", "支持：1–6人（按档）+样板间", "政策宣讲位"],
    ["会后逼单", "主责：7日回访报价跟进", "支持：看场会议室/决策人对接", "申报服务接棒"],
    ["成交交割", "协同：合同条款与服务包", "主责：租赁合同/交割/开票", "入驻后服务"],
    ["复盘", "主责：单场复盘+月报", "联席确认数据", "政策转化统计"],
]
table_block(s, Inches(1.2), rows, [0.14, 0.3, 0.32, 0.24], sizes=[11, 11, 11, 11], rh=0.52)

idx = nxt(); s = slide(); header(s, "02", "SUPPORT BY TIER", "分档现场人力（贵方支持上限清晰）"); footer(s, idx)
rows = [
    ["场次类型", "同浦汇", "园区支持", "服中心/其他", "合计约"],
    ["常规训练/沙龙（A/B）", "2–3人", "1–2人", "0–2人", "4–7人"],
    ["对接/撮合（C）", "3–4人", "2人", "0–1人", "5–7人"],
    ["内容市集/联展（D）", "4–5人", "3–4人+物业", "外包可选", "8–12人"],
    ["WAIC开放日（E）", "6人", "5人+物业", "服中心2+媒体", "15人+"],
    ["Builders/AGI夜（E）", "4–5人", "2–3人", "联合会/复旦2", "8–12人"],
    ["嘉年华级（E）", "6–8人", "5–6人+安保", "外包", "15人+"],
]
table_block(s, Inches(1.25), rows, [0.24, 0.16, 0.22, 0.22, 0.16], sizes=[11, 12, 12, 12, 12], rh=0.55)
text(s, ML, Inches(5.7), CW, Inches(0.9),
     "贵方最低支持清单（每场必做）：①场地报备 ②导视停车 ③样板间可看 ④在谈名单共享 ⑤会后1间逼单会议室。其余按档加配。", size=13, color=MUT)

# 收费
idx = nxt(); s = slide(); header(s, "03", "PRICING", "收费模式（活动费 + 佣金2个月 + 可选媒体）"); footer(s, idx)
card(s, ML, Inches(1.25), Inches(5.9), Inches(5.5), "活动打包怎么收费", accent=GOLD, items=[
    COMMERCIAL["activity_fee"],
    "大场（WAIC/嘉年华/联展）从包内高权重分摊",
    "培训营可另设企业付费席位（冲抵成本，规则书面约定）",
    "门票/赞助默认归园区经营性收入（可协商分成）",
    "媒体投流 " + COMMERCIAL["media_fee"].split("：")[1],
], bsize=12)
card(s, Inches(7.2), Inches(1.25), Inches(5.55), Inches(5.5), "招商怎么收费", accent=GREEN, items=[
    COMMERCIAL["commission"],
    COMMERCIAL["rent"],
    COMMERCIAL["property"],
    COMMERCIAL["rooms"],
    "活动转化成交：同样计佣金2个月，不对出租率对赌",
    "交付《意向→看场→报价→签约》漏斗周报",
], bsize=12)

# 交付物
idx = nxt(); s = slide(); header(s, "03", "DELIVERABLES", "我们具体交付什么（每场标准包）"); footer(s, idx)
rows = [
    ["交付物", "内容", "时点", "验收标准"],
    ["单场方案", "主题/议程/嘉宾/动线/预算分摊", "T-14", "双方书面确认"],
    ["邀约与关联度表", "名单+画像标签+负责人标记", "T-3", "邀约关联度≥90%"],
    ["现场执行", "签到建档/主持场控/物料", "D日", "流程按时完成"],
    ["意向分级表", "A/B/C级客户+跟进人", "T+1", "留资完整可回访"],
    ["7日回访纪要", "触达记录+看场预约+卡点", "T+7", "触达≥3次/线索"],
    ["单场复盘", "人数/关联度/负责人占比/转化", "T+14", "指标对照档位KPI"],
    ["月度联席包", "漏斗看板+下月排期", "每月", "联席会确认"],
]
table_block(s, Inches(1.2), rows, [0.18, 0.34, 0.14, 0.34], sizes=[11, 12, 11, 12], rh=0.55)

# 30场总表 概况（分两页）
for start, end in ((1, 15), (16, 30)):
    idx = nxt(); s = slide(); header(s, "04", "OVERVIEW", f"30场活动概况总表（{start}–{end}）"); footer(s, idx)
    rows = [["序号", "主题", "时间", "形式", "人数", "负责人%", "关联度", "场地"]]
    for e in EVENTS:
        if start <= e["id"] <= end:
            rows.append([str(e["id"]), e["name"][:14], e["month"], e["format"], e["people"], e["dm"], e["assoc"], e["floor"]])
    table_block(s, Inches(1.15), rows, [0.07, 0.28, 0.12, 0.12, 0.11, 0.1, 0.1, 0.1], sizes=[9, 10, 9, 9, 9, 9, 9, 9], rh=0.33)

# 价值+交付+分工+收费 总表（分两页，压缩列）
for start, end in ((1, 15), (16, 30)):
    idx = nxt(); s = slide(); header(s, "04", "VALUE & DELIVERY", f"价值 · 交付 · 分工 · 收费（{start}–{end}）"); footer(s, idx)
    rows = [["序", "活动价值（园区/招商）", "我方主责", "贵方支持", "收费"]]
    for e in EVENTS:
        if start <= e["id"] <= end:
            rows.append([
                str(e["id"]),
                e["value"][:22],
                "策划邀约执行回访",
                "场地名单看场1–N人",
                f"打包分摊~{e['budget_share']:.1f}万",
            ])
    table_block(s, Inches(1.15), rows, [0.06, 0.34, 0.22, 0.22, 0.16], sizes=[9, 10, 10, 10, 10], rh=0.33)

# 执行节奏
idx = nxt(); s = slide(); header(s, "04", "RHYTHM", "统一执行节奏（每场通用）"); footer(s, idx)
rows = [
    ["节点", "动作", "我方", "贵方", "产出"],
    ["T-21", "选题定档", "出方案草稿", "确认档期", "档期锁定"],
    ["T-14", "嘉宾物料", "邀约+物料定稿", "提供在谈名单", "方案确认"],
    ["T-7", "邀约确认", "确认到会与关联度", "领导出席答复", "到会名单"],
    ["T-1", "彩排", "流程彩排", "场地物业就绪", "可执行"],
    ["D日", "执行", "主持场控建档", "现场支持+看场", "签到+意向"],
    ["T+1", "名单", "分级建档", "共享CRM字段", "意向表"],
    ["T+7", "逼单", "回访报价", "会议室/决策人", "看场/报价"],
    ["T+14", "复盘", "复盘报告", "联席确认", "月报入库"],
]
table_block(s, Inches(1.2), rows, [0.1, 0.16, 0.24, 0.24, 0.26], sizes=[11, 12, 12, 12, 12], rh=0.48)

# 案例页
for ci, case in enumerate(CASES):
    idx = nxt(); s = slide(); header(s, "05", f"CASE {ci+1}", f"往期类似案例｜{case['name']}"); footer(s, idx)
    text(s, ML, Inches(1.15), CW, Inches(0.35),
         [(f"{case['when']}　｜　{case['where']}　｜　{case['tier']}", GOLD, True)], size=12)
    card(s, ML, Inches(1.55), Inches(6.3), Inches(2.55), "可见机构/企业", accent=GOLD, items=[
        case["orgs"][:60], case["companies"][:70], f"来源：{case['source']}", case["map_to"],
    ], bsize=11)
    card(s, Inches(7.05), Inches(1.55), Inches(5.7), Inches(2.55), "当天议程", accent=ACC, items=case["agenda"], bsize=11)
    photos = case["photos"][:3]
    pw = Inches(3.85); ph = Inches(2.35); gap = Inches(0.2)
    for i, fn in enumerate(photos):
        x = Emu(ML + i * (pw + gap))
        add_pic(s, os.path.join(PHOTO, fn), x, Inches(4.3), pw, ph)
    text(s, ML, Inches(6.72), CW, Inches(0.25), "现场照片来源：杨浦区政府官网/公开报道配图（用于同类活动能力证明）", size=9, color=SOFT)

# 案例映射
idx = nxt(); s = slide(); header(s, "05", "CASE MAP", "各档活动 ↔ 往期案例映射"); footer(s, idx)
rows = [
    ["创智汇档位", "对应往期案例", "可迁移的执行要点"],
    ["A 训练营/黑客松", "融见科创路演（项目展示段）", "路演评审机制+当场看场"],
    ["B 政策/治理沙龙", "杨数浦AI/合规沙龙", "政策解读+案例+一对一咨询"],
    ["C 产业对接", "融见科创6+6+6+1对接", "供需桌+金融机构+政府代表"],
    ["D 内容活动", "创作者/市集级（待补同浦汇实拍）", "发布+快闪+招商通道"],
    ["E WAIC大场", "滨江授证仪式级现场", "开放参观+仪式+集中洽谈"],
    ["F 出海/领事", "杨数浦出海合规沙龙", "合规顾问+圆桌+授证背书"],
]
table_block(s, Inches(1.25), rows, [0.22, 0.36, 0.42], sizes=[12, 12, 12], rh=0.58)
text(s, ML, Inches(5.9), CW, Inches(0.7),
     "说明：公开案例证明「我们能办到什么规格」；创智汇落地后，每档将沉淀本场地实拍与企业名单，替换进月度案例库。", size=12, color=MUT)

# 结论
idx = nxt(); s = slide(); header(s, "END", "SUMMARY", "一页纸结论（可直接回复甲方三点）"); footer(s, idx)
card(s, ML, Inches(1.25), Inches(4.0), Inches(5.5), "①关联度", accent=GOLD, items=[
    "邀约≥90% / 到场≥80%", "负责人占比按档35–60%", "不是口号100%，可书面锁名单", "每场交付关联度核验表",
], bsize=12)
card(s, Inches(4.75), Inches(1.25), Inches(4.0), Inches(5.5), "②分工", accent=ACC, items=[
    "我方：策划邀约执行回访复盘", "贵方：场地物业名单看场交割", "最低支持清单5项写死", "分档人力表已列明",
], bsize=12)
card(s, Inches(8.9), Inches(1.25), Inches(3.85), Inches(5.5), "③案例+收费", accent=GREEN, items=[
    "3个公开案例含企业/议程/照片", "活动30场/30万打包", "佣金2个月·租金3.3·物业13.8", "8间待租优先活动转化",
], bsize=12)

total = len(_FOOT)
for run in _FOOT:
    run.text = " / %02d" % total

os.makedirs(ART, exist_ok=True)
outs = [
    ("创智汇30场活动具体方案.pptx", "chuangzhihui-30-events-plan.pptx"),
    ("创智汇30场活动具体方案-WAIC对齐版.pptx", "chuangzhihui-30-events-waic-aligned.pptx"),
    ("创智汇30场活动具体方案-表格应答版.pptx", "chuangzhihui-30-events-table-v3.pptx"),
]
tmp = os.path.join(HERE, "_tmp_events_v3.pptx")
prs.save(tmp)
for cn, en in outs:
    p1 = os.path.join(HERE, cn); p2 = os.path.join(HERE, en)
    shutil.copy2(tmp, p1); shutil.copy2(tmp, p2)
    shutil.copy2(p1, os.path.join(ART, cn)); shutil.copy2(p2, os.path.join(ART, en))
os.remove(tmp)
print("PPT slides:", len(prs.slides))

# ===================== EXCEL =====================
wb = Workbook()
thin = Border(
    left=Side(style="thin", color="D0D0D0"), right=Side(style="thin", color="D0D0D0"),
    top=Side(style="thin", color="D0D0D0"), bottom=Side(style="thin", color="D0D0D0"),
)
head_fill = PatternFill("solid", fgColor="4A2C7A")
alt_fill = PatternFill("solid", fgColor="F3E9FF")
warn_fill = PatternFill("solid", fgColor="FFF2CC")


def style_header(ws, headers):
    for j, h in enumerate(headers, 1):
        cell = ws.cell(1, j, h)
        cell.font = Font(name="微软雅黑", bold=True, color="FFFFFF", size=10)
        cell.fill = head_fill
        cell.alignment = Alignment(wrap_text=True, vertical="center", horizontal="center")
        cell.border = thin


def write_rows(ws, rows, widths, row_h=42):
    for i, row in enumerate(rows, 2):
        for j, v in enumerate(row, 1):
            cell = ws.cell(i, j, v)
            cell.font = Font(name="微软雅黑", size=9)
            cell.alignment = Alignment(wrap_text=True, vertical="center")
            cell.border = thin
            if i % 2 == 0:
                cell.fill = alt_fill
        ws.row_dimensions[i].height = row_h
    ws.row_dimensions[1].height = 30
    for i, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = w


# Sheet1 概况
ws = wb.active; ws.title = "01-活动概况总表"
headers = ["序号", "档位", "WAIC赛道", "主题/活动名称", "时间", "活动形式", "执行节奏", "人数",
           "适配企业负责人占比", "孵化器关联度目标", "场地", "时长"]
style_header(ws, headers)
rows = []
for e in EVENTS:
    rows.append([
        e["id"], CAT_META[e["cat"]][0], e["waic"], e["name"], e["month"], e["format"], e["rhythm"],
        e["people"], e["dm"], e["assoc"], e["floor"], e["dur"],
    ])
write_rows(ws, rows, [6, 14, 12, 28, 12, 10, 36, 10, 12, 12, 14, 12], 48)

# Sheet2 价值交付分工收费
ws2 = wb.create_sheet("02-价值交付分工收费")
headers2 = ["序号", "活动名称", "活动价值", "具体交付物", "我方（同浦汇）做什么", "贵方（园区）支持什么",
            "联合会/服中心", "收费模式", "打包分摊约(万)"]
style_header(ws2, headers2)
rows2 = []
for e in EVENTS:
    rows2.append([
        e["id"], e["name"], e["value"], e["deliver"], e["us"], e["park"], e["union"], e["fee"], e["budget_share"],
    ])
write_rows(ws2, rows2, [6, 26, 28, 40, 28, 28, 22, 36, 10], 56)

# Sheet3 关联度口径
ws3 = wb.create_sheet("03-孵化器关联度口径")
style_header(ws3, ["问题", "明确答复", "可执行标准"])
rows3 = [
    ["是不是100%与孵化器关联？", "不是口号式“到场100%都是入孵企业”", ASSOC_RULE["not_100"]],
    ["邀约标准", ASSOC_RULE["invite_target"], "报名表行业标签审核，不达标不发正式邀请函"],
    ["到场标准", ASSOC_RULE["onsite_target"], "签到核验；大场≥75%（D/E流量场）"],
    ["负责人占比", ASSOC_RULE["dm_note"], "名片/职位字段统计，写入单场复盘"],
    ["若甲方要求更严", ASSOC_RULE["if_require_100"], "按场次书面锁定邀约名单审核权"],
]
write_rows(ws3, rows3, [28, 40, 50], 60)
for r in range(2, 7):
    ws3.cell(r, 1).fill = warn_fill

# Sheet4 分工RACI
ws4 = wb.create_sheet("04-执行分工RACI")
style_header(ws4, ["环节", "同浦汇（我方）", "园区/科创集团（贵方）", "联合会/服中心", "成功标准"])
rows4 = [
    ["选题定档", "主责年度日历与单场方案", "确认档期与禁办日", "政策窗口协同", "档期锁定"],
    ["嘉宾邀约", "主责企业/导师/投资人邀约", "共享在谈客户名单", "会员协同邀约", "邀约关联度≥90%"],
    ["场地物料", "活动物料/签到/主持", "场地物业水电导视停车", "领导出席协调", "T-1就绪"],
    ["现场执行", "流程场控建档", "1–6人支持+样板间开放", "政策宣讲位", "留资可回访"],
    ["会后逼单", "7日回访报价跟进", "会议室/决策人对接", "申报服务接棒", "看场率≥25%"],
    ["成交交割", "协同合同与服务包", "租赁合同/交割/开票", "入驻后服务", "签约入库"],
    ["复盘", "单场复盘+月报", "联席确认数据", "政策转化统计", "KPI对照达成"],
]
write_rows(ws4, rows4, [12, 28, 28, 22, 18], 40)

# Sheet5 收费与租金
ws5 = wb.create_sheet("05-收费与租金佣金")
style_header(ws5, ["条款", "标准", "备注"])
rows5 = [
    ["办公租金", "3.3 元/㎡/天", "3F办公/孵化单元"],
    ["物业费", "13.8 元/㎡/月", "按月计收"],
    ["待租办公室", "8 间", "活动优先导流转化"],
    ["招商佣金", "成交首年有效净租金 × 2 个月", "首年不重复；活动转化同样适用"],
    ["活动执行费", "30 场 / 年 打包 30 万元", "含策划执行、基础物料、主持协调"],
    ["媒体流量费", "约 10 万元/年（可选）", "另计"],
    ["门票/赞助", "默认归园区经营性收入", "可协商分成"],
    ["不对赌", "不对出租率/去化面积对赌", "轻资产联合运营"],
]
write_rows(ws5, rows5, [16, 36, 40], 32)

# Sheet6 往期案例
ws6 = wb.create_sheet("06-往期类似案例")
style_header(ws6, ["档位映射", "案例名称", "时间", "地点", "可见机构/企业", "当天议程", "现场照片文件", "来源", "映射创智汇场次"])
rows6 = []
for c in CASES:
    rows6.append([
        c["tier"], c["name"], c["when"], c["where"], f"{c['orgs']}；企业侧：{c['companies']}",
        " | ".join(c["agenda"]), "；".join(c["photos"]), c["source"], c["map_to"],
    ])
write_rows(ws6, rows6, [14, 36, 14, 28, 40, 50, 36, 28, 28], 70)

# Sheet7 行程明细
ws7 = wb.create_sheet("07-逐场行程嘉宾")
style_header(ws7, ["序号", "活动名称", "嘉宾构成", "行程", "园区价值/招商价值"])
rows7 = [[e["id"], e["name"], e["guests"], e["agenda"], e["value"]] for e in EVENTS]
write_rows(ws7, rows7, [6, 28, 36, 50, 36], 50)

x_files = [
    "创智汇30场活动方案排期表.xlsx",
    "chuangzhihui-30-events-schedule.xlsx",
    "创智汇30场活动方案排期表-WAIC对齐.xlsx",
    "chuangzhihui-30-events-schedule-waic.xlsx",
    "创智汇30场活动方案-表格应答版.xlsx",
    "chuangzhihui-30-events-table-v3.xlsx",
]
for fn in x_files:
    p = os.path.join(HERE, fn)
    wb.save(p)
    shutil.copy2(p, os.path.join(ART, fn))
print("Excel OK")
