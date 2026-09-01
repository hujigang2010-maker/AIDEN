# -*- coding: utf-8 -*-
"""创智汇 · 30场活动专项交付（双会资源融合版）
口径沿用修订版；本版重点：
- 结合 WAIC2026 全量资源整合总表 + ChinaJoy2026 招商引资总表
- 只取与招商贴近的层次与叠加，不做全量名单堆砌
- 为30场做「定义核 + 资源导入」映射，响应「AI+数字内容无界共创港」
"""
import os, shutil
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

HERE = os.path.dirname(os.path.abspath(__file__))
ART = "/opt/cursor/artifacts"
PHOTO = os.path.join(HERE, "cases", "photos")
FONT = "Microsoft YaHei"

INK = RGBColor(0x1A, 0x1A, 0x2E)
MUT = RGBColor(0x4A, 0x4A, 0x5A)
SOFT = RGBColor(0x8A, 0x8A, 0x9A)
BLUE = RGBColor(0x1E, 0x4D, 0x8C)
GOLD = RGBColor(0xC4, 0x9A, 0x3C)
LINE = RGBColor(0xD8, 0xDE, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF7, 0xF8, 0xFA)
ACC = RGBColor(0xE8, 0xF0, 0xFA)
GREEN = RGBColor(0x1F, 0x7A, 0x5C)
ROSE = RGBColor(0xA6, 0x3D, 0x40)

SW, SH = Inches(13.333), Inches(7.5)
ML, CW = Inches(0.55), Inches(12.23)

PARTIES = {
    "ops": "上海市云计算创新基地（国家级孵化器）——纳入主体运营范围",
    "academic": "复旦大学住房政策研究中心——学术支持单位",
    "support": "上海市杨浦区科技企业联合会、上海市杨浦区科技企业服务中心——活动及运营载体支持单位",
    "exec": "同浦汇——活动策划执行与带客",
}

# 双会资源池（只取招商可用摘要，不落全量）
WAIC_POOL = {
    "scale": "参展商963 · 品牌主库4262 · 联系人3463 · A/B线索500 · 论坛175",
    "tracks": ["算力·芯片·基础设施", "机器人与具身", "大模型/Agent", "AIGC内容创意", "行业AI", "治理合规"],
    "act": [
        ("ACT-001", "大模型/Agent闭门", "A2/A3/F3"),
        ("ACT-002", "机器人体验+供需", "A5/E1"),
        ("ACT-003", "算力芯片小型峰会", "A4/B4"),
        ("ACT-006", "行业AI案例圆桌", "C4/B1"),
    ],
}

CJ_POOL = {
    "theme": "与AI同游",
    "scale": "近900家参展 · 可触达优先263 · 精编422",
    "priority": [
        ("最高", "中小游戏工作室池", "Express/Next Play/AGS → C1/E4/D2"),
        ("高", "具身智能 Vision Future", "宇树/逐际类线索 → A5/E1"),
        ("高", "潮玩谷子供应链", "魔玩/创作者 → C3/D3/C5"),
        ("中高", "创作者经济", "创作者展区 → D2/D5"),
    ],
    "fit": "中腰部工作室选址灵活；不硬追大厂总部搬迁",
}

# 30场：人次≤30；负责人≤30%；每场挂资源导入定义核
EVENTS = [
    dict(code="P1", cat="P", name="创智汇项目推广日①·中介与渠道专场", month="2026年8月上旬", fmt="项目推广日",
         people="≤30人", dm="≤30%", floor="5F+3F", rhythm="T-14物料→T-7邀约→T-1彩排→D日→建群回访",
         lease="中介渠道带看与项目推介",
         agenda="项目介绍→空间导览→政策要点→一对一洽谈→建群",
         guests="房产/产业中介、渠道机构、园区招商同事", value="渠道点火；待租单元曝光",
         deliver="方案+议程+签到+意向表+建群+回访摘要+复盘",
         layer="L0点火", src="园区自建", import_res="双会客群转介话术包；待租8间看场清单",
         define="渠道点火场：把「AI+数字内容共创港」讲清楚，导入后续带客池"),
    dict(code="P2", cat="P", name="创智汇项目推广日②·五大行与金融机构专场", month="2026年8月中旬", fmt="项目推广日",
         people="≤30人", dm="≤30%", floor="5F沙龙", rhythm="同上",
         lease="金融资源链接与企业客户导入",
         agenda="项目路演→金融产品对接→企业需求座谈→看场→建群",
         guests="五大行及合作金融机构对公/科创条线", value="金融背书；客户转介",
         deliver="同上", layer="L0点火", src="园区自建+投研交叉",
         import_res="WAIC金融科技论坛议题摘要；CJ BTOB专业观众画像",
         define="金融转介场：用科创金融话术放大带客可信度"),
    dict(code="P3", cat="P", name="创智汇项目推广日③·投研机构专场", month="2026年8月中旬", fmt="项目推广日",
         people="≤30人", dm="≤30%", floor="5F圆桌", rhythm="同上",
         lease="投研视角放大项目影响力",
         agenda="项目定位→产业赛道研判→圆桌→看场→建群",
         guests="券商研究所、产业研究院、智库分析师", value="专业传播；精准客群",
         deliver="同上", layer="L0点火", src="园区自建",
         import_res="WAIC赛道结论（大模型/具身/算力）；CJ「与AI同游」产业判断",
         define="叙事定调场：对外统一园区赛道故事"),
    dict(code="P4", cat="P", name="创智汇项目推广日④·政府部门与载体协同专场", month="2026年8月下旬", fmt="项目推广日",
         people="≤30人", dm="≤30%", floor="3F+5F", rhythm="同上",
         lease="政策与载体协同落地",
         agenda="载体介绍→政策协同→云创基地资源→看场→建群",
         guests="区相关部门、载体代表、云创基地、服中心", value="政策协同；主体背书",
         deliver="同上", layer="L0点火", src="云创基地+科企联/服中心",
         import_res="云创算力券/创新券口径；YOUNG立方政策包",
         define="政策协同场：主体背书与载体资源上桌"),
    dict(code="E1", cat="E", name="WAIC成果承接·创智汇开放交流日", month="2026年8月下旬", fmt="开放交流日",
         people="≤30人", dm="≤30%", floor="3F+5F", rhythm="同上",
         lease="衔接WAIC议题与园区看场",
         agenda="WAIC议题速递→园区承接→双展台→洽谈→建群",
         guests="AI/内容企业、媒体、云创基地代表", value="WAIC衔接；看场转化",
         deliver="同上+通稿", layer="L1双会承接", src="WAIC总表",
         import_res="A/B线索抽样邀约；ACT-001/002/006议题切片；不堆全量963展商",
         define="WAIC余热承接：把大会议题压缩为园区可看场的≤30人交流"),
    dict(code="A2", cat="A", name="Agent智能体搭建一日营", month="2026年9月上旬", fmt="实训营",
         people="≤30人", dm="≤30%", floor="3F OPC", rhythm="同上",
         lease="Agent初创落位", agenda="签到→工作流→实操→Demo→看场",
         guests="Agent工程师、初创团队", value="入孵转化", deliver="同上",
         layer="L2·3F能力", src="WAIC·大模型/Agent",
         import_res="WAIC行业：AI Agent约87家池；ACT-001闭门会缩尺",
         define="3F入孵核：Agent能力→OPC工位转化"),
    dict(code="C1", cat="C", name="ChinaJoy主题·数字娱乐与IP授权对接会", month="2026年9月中旬", fmt="对接会",
         people="≤30人", dm="≤30%", floor="5F展示中心", rhythm="同上",
         lease="ChinaJoy赛道企业/IP方入驻与展位",
         agenda="ChinaJoy赛道解读→IP授权→一对一撮合→看场→建群",
         guests="游戏/动漫/潮玩IP方、渠道、内容公司", value="衔接ChinaJoy；展位/办公转化",
         deliver="同上", layer="L2·5F内容", src="CJ·游戏风云/IP",
         import_res="CJ优先：中小工作室+IP授权方；可触达263抽样，不追大厂总部",
         define="5F内容核：IP授权对接→展位/办公问询"),
    dict(code="F1", cat="F", name="出海专题活动（另议·另计价）", month="档期另议", fmt="专题（另议）",
         people="≤30人", dm="≤30%", floor="另议", rhythm="立项后单独排期",
         lease="出海需求企业服务（不在年度活动包费用内）",
         agenda="另案确定", guests="另案确定",
         value="出海场次价格与方案另议，不计入本册活动服务费",
         deliver="另案合同与交付清单", layer="另计价", src="CJ·国际展团（另案）",
         import_res="CJ BTOB外资线索仅作另案素材；本册不承诺外企到场",
         define="另议核：出海不进年度包，另案另计价"),
    dict(code="A3", cat="A", name="多模态智能体工作坊", month="2026年10月上旬", fmt="工作坊",
         people="≤30人", dm="≤30%", floor="3F培训", rhythm="同上",
         lease="多模态团队入驻", agenda="案例→实操→共创→评审→看场",
         guests="多模态应用团队、产品负责人", value="技术团队入驻", deliver="同上",
         layer="L2·3F能力", src="WAIC·多模态/AIGC",
         import_res="WAIC AIGC图像/视频/3D子类抽样；论坛「内容创意与AIGC」切片",
         define="3F技术核：多模态团队看场入驻"),
    dict(code="E4", cat="E", name="ChinaJoy主题日·数字娱乐生态交流", month="2026年10月中旬", fmt="主题日",
         people="≤30人", dm="≤30%", floor="5F", rhythm="同上",
         lease="ChinaJoy生态企业集中看场与洽谈",
         agenda="开幕→赛道分享→展洽→看场→建群",
         guests="数字娱乐企业、内容工作室、渠道", value="衔接ChinaJoy；集中洽谈",
         deliver="同上", layer="L2·5F内容", src="CJ·游戏风云",
         import_res="Express/独立游戏中小团队池；库洛/鹰角/散爆等仅作赛道示例非名单承诺",
         define="5F生态核：数字娱乐集中看场"),
    dict(code="B2", cat="B", name="AI治理与可信智能体沙龙", month="2026年10月下旬", fmt="沙龙",
         people="≤30人", dm="≤30%", floor="3F", rhythm="同上",
         lease="合规型企业信任导入", agenda="治理框架→合规清单→案例→问答→对接",
         guests="合规顾问、企业法务/负责人", value="信任入驻", deliver="同上",
         layer="L2·3F能力", src="WAIC·治理标准论坛",
         import_res="WAIC治理赛道7场论坛议题摘要；合规顾问短名单",
         define="信任核：治理合规降低入驻决策摩擦"),
    dict(code="A4", cat="A", name="火山引擎×算力Infra实务营", month="2026年11月上旬", fmt="厂商联训",
         people="≤30人", dm="≤30%", floor="3F+云创基地", rhythm="同上",
         lease="高算力企业定向看场", agenda="政策包→代金券→案例→诊断→看场",
         guests="火山引擎、云创基地、企业技术负责人", value="云创协同；算力企业", deliver="同上",
         layer="L2·3F能力", src="WAIC·算力/芯片 + CJ云服务示例",
         import_res="ACT-003缩尺；WAIC算力展商226家池抽样；云创算力资源",
         define="算力核：云创基地×厂商联训→高算力看场"),
    dict(code="B4", cat="B", name="创新券·算力券·模型券实务沙龙", month="2026年11月中旬", fmt="实务沙龙",
         people="≤30人", dm="≤30%", floor="3F", rhythm="同上",
         lease="用券企业聚集", agenda="三券规则→核销→案例→开户→入驻洽谈",
         guests="券平台、云厂商、企业负责人", value="券务转化", deliver="同上",
         layer="L2·3F能力", src="云创+WAIC算力池",
         import_res="三券核销路径；WAIC算力/云厂商联系人抽样",
         define="券务核：用券企业聚集转化"),
    dict(code="L1", cat="L", name="领事到访接待（另计价）", month="档期另议", fmt="外事接待（另计价）",
         people="另议", dm="—", floor="会客厅/5F", rhythm="外事流程单独排期",
         lease="领事到访接待（不在三部分活动费用内）",
         agenda="另案确定", guests="领事及相关方（另案）",
         value="领事到访与挂牌活动分属不同性质，均另计价",
         deliver="另案报价与执行方案", layer="另计价", src="另案",
         import_res="不占用年度30场资源池",
         define="外事另计价核：到访≠挂牌，分开报价"),
    dict(code="A5", cat="A", name="具身智能空间交互体验日", month="2026年12月上旬", fmt="体验日",
         people="≤30人", dm="≤30%", floor="5F+3F", rhythm="同上",
         lease="具身/机器人团队看场", agenda="演示→讲解→场景→踩点→洽谈",
         guests="具身团队、高校实验室", value="具身看场", deliver="同上",
         layer="L3叠加", src="WAIC具身 + CJ Vision Future",
         import_res="ACT-002缩尺；WAIC场景观察（零售/工业物流）；CJ前沿科技高优先级",
         define="双会叠加核：具身体验→3F/5F联合看场"),
    dict(code="F4", cat="F", name="创智汇AI年度Demo Day", month="2026年12月中旬", fmt="路演日",
         people="≤30人", dm="≤30%", floor="5F主展", rhythm="同上",
         lease="集中签约洽谈与媒体背书", agenda="开幕→精选路演→颁奖→洽谈",
         guests="投资人、链主、媒体、精选项目", value="集中签约", deliver="同上+签约台账",
         layer="L4收口", src="双会精选项目",
         import_res="年内WAIC/CJ导入的A/B意向项目精选路演，不做全量路演",
         define="年中收口核：精选项目集中签约洽谈"),
    dict(code="D1", cat="D", name="AIGC微短剧制片特训", month="2027年1月上旬", fmt="特训营",
         people="≤30人", dm="≤30%", floor="3F+5F", rhythm="同上",
         lease="厂牌/工作室入驻", agenda="政策→制片→脚本→路演→看场",
         guests="导演、厂牌、平台方", value="内容入驻", deliver="同上",
         layer="L2·5F内容", src="WAIC·内容创意AIGC",
         import_res="WAIC「AI赋能影视内容」论坛切片；AIGC厂牌抽样",
         define="内容制片核：微短剧厂牌→5F展陈/3F办公"),
    dict(code="A6", cat="A", name="AI营销Agent实战营", month="2027年1月中旬", fmt="实战营",
         people="≤30人", dm="≤30%", floor="3F", rhythm="同上",
         lease="营销科技公司看场", agenda="策略→Agent→素材→复盘→转化",
         guests="投放操盘手、产品经理", value="营销科技入驻", deliver="同上",
         layer="L2·3F能力", src="WAIC·企业服务与营销",
         import_res="WAIC企业服务/营销86家池抽样；Agent应用团队",
         define="营销科技核：投放团队看场入驻"),
    dict(code="B1", cat="B", name="YOUNG立方×智能伙伴政策沙龙", month="2027年2月中旬", fmt="政策沙龙",
         people="≤30人", dm="≤30%", floor="5F沙龙", rhythm="同上",
         lease="内容/AI企业导入", agenda="政策→礼包→画像→诊断→看场",
         guests="政策宣讲、服中心、企业负责人", value="政策获客", deliver="同上",
         layer="L2·政策", src="服中心+双会客群",
         import_res="政策礼包+双会导入意向企业复邀",
         define="政策获客核：YOUNG立方礼包转化"),
    dict(code="B3", cat="B", name="高企认定冲刺（AI企业专场）", month="2027年2月下旬", fmt="辅导会",
         people="≤30人", dm="≤30%", floor="3F", rhythm="同上",
         lease="待认定企业带政策入驻", agenda="条件→材料→初筛→套餐→激励",
         guests="辅导顾问、待认定企业负责人", value="带政策入驻", deliver="同上",
         layer="L2·政策", src="服中心",
         import_res="双会AI应用类待认定企业名单抽样",
         define="带政策入驻核：高企辅导绑定入驻套餐"),
    dict(code="A7", cat="A", name="OPC超级个体黑客松（春）", month="2027年3月上旬", fmt="黑客松",
         people="≤30人", dm="≤30%", floor="3F", rhythm="同上",
         lease="获奖团队优先谈单元", agenda="开题→开发→路演→礼包→看场",
         guests="评委、投资人、Builders", value="获奖谈单元", deliver="同上",
         layer="L2·3F能力", src="WAIC·青年/OPC + CJ创作者",
         import_res="高校/青年Builder；CJ AGS高校团队线索抽样",
         define="OPC核：黑客松获奖→优先谈单元"),
    dict(code="C2", cat="C", name="高校成果转化·AI for Science日", month="2027年3月下旬", fmt="对接日",
         people="≤30人", dm="≤30%", floor="3F+沙龙", rhythm="同上",
         lease="成果公司/实验室落户", agenda="成果路演→场景→承接→看单元→洽谈",
         guests="复旦/同济技转、教授团队、住房政策研究中心学者（学术支持）", value="成果落户", deliver="同上",
         layer="L2·学术", src="学术支持+WAIC科研赛道",
         import_res="复旦住房政策研究中心学术支持位；WAIC前沿科研论坛切片",
         define="成果转化核：高校实验室/公司落户"),
    dict(code="F3", cat="F", name="通往AGI季度圆桌", month="2027年3月下旬", fmt="闭门圆桌",
         people="≤30人", dm="≤30%", floor="5F沙龙", rhythm="同上",
         lease="研究型/模型团队", agenda="议题→圆桌→围炉→问答→看场",
         guests="学者、模型企业、投资人", value="研究型客户", deliver="同上",
         layer="L2·3F能力", src="WAIC·大模型论坛",
         import_res="ACT-001闭门形态；模型公司抽样（不作全量邀约）",
         define="研究型客户核：闭门圆桌深谈看场"),
    dict(code="D2", cat="D", name="ChinaJoy主题·创作者内容首发①", month="2027年4月中旬", fmt="发布会",
         people="≤30人", dm="≤30%", floor="5F", rhythm="同上",
         lease="数字娱乐/内容品牌问询", agenda="揭幕→发布→快闪→专访→招商通道",
         guests="创作者、媒体、渠道、ChinaJoy相关内容方", value="衔接ChinaJoy内容侧", deliver="同上",
         layer="L2·5F内容", src="CJ·创作者展区",
         import_res="CJ创作者/社团抽样；渠道复用（抖音/小红书玩法迁移园区）",
         define="创作者首发核：内容发布带招商通道"),
    dict(code="C3", cat="C", name="ChinaJoy主题·游戏周边与潮玩供应链对接", month="2027年4月下旬", fmt="对接会",
         people="≤30人", dm="≤30%", floor="5F", rhythm="同上",
         lease="ChinaJoy供应链/周边企业展位落位",
         agenda="赛道介绍→供需对接→展位参观→报价→建群",
         guests="游戏周边、潮玩供应链、渠道采购", value="ChinaJoy供应链去化", deliver="同上",
         layer="L2·5F内容", src="CJ·魔玩/谷子",
         import_res="CJ潮玩谷子高优先级；JOYTOY/卡游等仅作品类示例",
         define="供应链核：周边/潮玩→5F展位落位"),
    dict(code="C4", cat="C", name="专精特新·AI应用培育路演", month="2027年5月中旬", fmt="路演",
         people="≤30人", dm="≤30%", floor="5F", rhythm="同上",
         lease="成长型AI补位", agenda="政策→路演→点评→金融→看场",
         guests="顾问、银行、基金、企业", value="成长型企业", deliver="同上",
         layer="L2·3F能力", src="WAIC·行业AI + ACT-006",
         import_res="ACT-006案例圆桌缩尺；行业AI114家池抽样",
         define="成长补位核：专精特新AI应用看场"),
    dict(code="D3", cat="D", name="ChinaJoy主题·数字娱乐市集体验日", month="2027年5月下旬", fmt="体验日",
         people="≤30人", dm="≤30%", floor="5F", rhythm="同上",
         lease="优质摊主/工作室升级固定展位",
         agenda="布展→体验→交流→转正洽谈→建群",
         guests="内容工作室、周边品牌、达人", value="ChinaJoy氛围转化", deliver="同上",
         layer="L2·5F内容", src="CJ·创作者/市集形态",
         import_res="CJ创作者展区小微品牌；市集→固定展位升级路径",
         define="市集转化核：摊主升级固定展位"),
    dict(code="C5", cat="C", name="ChinaJoy主题·IP联名与衍生品撮合会", month="2027年6月中旬", fmt="撮合会",
         people="≤30人", dm="≤30%", floor="5F", rhythm="同上",
         lease="IP联名/衍生品团队入驻",
         agenda="IP路演→联名模式→一对一→MOU→看场",
         guests="IP方、衍生品商、渠道", value="ChinaJoy IP商业化落地", deliver="同上",
         layer="L2·5F内容", src="CJ·IP/联名",
         import_res="CJ老字号×IP联名案例迁移；衍生品商抽样",
         define="IP商业化核：联名撮合→入驻/展位"),
    dict(code="D4", cat="D", name="ChinaJoy主题·沉浸式数字娱乐联展", month="2027年6月下旬", fmt="联展（短展期）",
         people="≤30人/场次接待", dm="≤30%", floor="5F", rhythm="展期预约制",
         lease="闭幕洽谈集中转化", agenda="布展→预约观展→交流→闭幕洽谈",
         guests="联合内容方、渠道、品牌", value="ChinaJoy视觉落地", deliver="同上",
         layer="L2·5F内容", src="CJ·氛围落地",
         import_res="联合内容方短名单；预约制控流≤30人/场",
         define="视觉落地核：短展期氛围→闭幕洽谈"),
    dict(code="D5", cat="D", name="创作者首发②·衔接WAIC2027与ChinaJoy", month="2027年7月中旬", fmt="发布会",
         people="≤30人", dm="≤30%", floor="5F", rhythm="同上",
         lease="旺季补位，双大会预热", agenda="发布→渠道→看场→年框→双会预热",
         guests="IP方、渠道、媒体", value="同步衔接WAIC与ChinaJoy", deliver="同上",
         layer="L4收口/预热", src="双会预热",
         import_res="全年双会导入线索复盘；预热WAIC2027+下届CJ",
         define="双会预热核：创作者首发衔接下一年双会"),
]

assert len(EVENTS) == 30

CASES = [
    {
        "name": "杨「数」浦数字沙龙：AI如何重塑企业DNA",
        "meta": "2025-06-30 · 美团上海综合指挥中心 · 沙龙类",
        "orgs": "杨浦区委网信办；赛博院等",
        "companies": "美团；大众点评大模型团队等",
        "agenda": ["参观数字化场景", "主题分享", "案例落地", "政策解读", "互动答疑"],
        "photos": ["case_c987bd91.jpg", "case_59157ea5.jpg", "case_6e51af9e.jpg"],
        "map": "对应：B类政策沙龙、A类训练营",
    },
    {
        "name": "「融见科创」人工智能专场路演暨投融资对接会",
        "meta": "2025-10 · 杨浦 · 路演对接类",
        "orgs": "杨浦科创促进会、邮储银行等",
        "companies": "复楚智能、卡房信息、一造科技、中科趋势、万笔千墨等",
        "agenda": ["开场致辞", "主办致辞", "趋势分享", "项目路演对接", "点评跟进"],
        "photos": ["case_db1d5cac.jpg", "case_65931516.jpg", "case_cb0a3812.jpg"],
        "map": "对应：F4 Demo Day、项目推广日、路演类",
    },
]


def rect(s, x, y, w, h, fill=None, line=None, lw=0.8):
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    b.shadow.inherit = False
    if fill is None: b.fill.background()
    else: b.fill.solid(); b.fill.fore_color.rgb = fill
    if line is None: b.line.fill.background()
    else: b.line.color.rgb = line; b.line.width = Pt(lw)
    return b


def text(s, x, y, w, h, content, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = content; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = FONT
    return tb


def bullets(s, x, y, w, h, items, size=12, color=MUT, mark=BLUE):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4); p.line_spacing = 1.15
        r0 = p.add_run(); r0.text = "▪ "; r0.font.size = Pt(size); r0.font.color.rgb = mark; r0.font.name = FONT
        r = p.add_run(); r.text = it; r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = FONT
    return tb


def slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, SW, SH, fill=WHITE)
    rect(s, 0, 0, SW, Inches(0.06), fill=BLUE)
    return s


def header(s, sec, title, sub=None):
    text(s, ML, Inches(0.22), Inches(12), Inches(0.28), sec, size=12, color=BLUE, bold=True)
    text(s, ML, Inches(0.48), Inches(12), Inches(0.42), title, size=22, color=INK, bold=True)
    if sub:
        text(s, ML, Inches(0.92), Inches(12), Inches(0.28), sub, size=12, color=MUT)


def footer(s, page, total=24):
    text(s, ML, Inches(7.15), Inches(9.2), Inches(0.25),
         "上海创智汇 × 同浦汇 · 活动专项 · 双会资源导入（WAIC & ChinaJoy）", size=10, color=SOFT)
    text(s, Inches(11.3), Inches(7.15), Inches(1.5), Inches(0.25), f"{page}/{total}", size=10, color=SOFT, align=PP_ALIGN.RIGHT)


def card(s, x, y, w, h, title, items, accent=BLUE):
    rect(s, x, y, w, h, fill=BG, line=LINE)
    rect(s, x, y, Inches(0.08), h, fill=accent)
    text(s, Emu(x + Inches(0.22)), Emu(y + Inches(0.12)), Emu(w - Inches(0.35)), Inches(0.32), title, size=14, color=INK, bold=True)
    bullets(s, Emu(x + Inches(0.22)), Emu(y + Inches(0.48)), Emu(w - Inches(0.4)), Emu(h - Inches(0.55)), items, size=12, mark=accent)


def table(s, y0, rows, col_w, sizes=None, rh=0.32):
    yy = y0
    if sizes is None: sizes = [10] * len(rows[0])
    for ri, row in enumerate(rows):
        xx = ML; h = Inches(0.38 if ri == 0 else rh)
        for ci, val in enumerate(row):
            cw = Emu(int(CW * col_w[ci]))
            bg = BLUE if ri == 0 else (ACC if ri % 2 == 0 else WHITE)
            fg = WHITE if ri == 0 else INK
            rect(s, xx, yy, cw, h, fill=bg, line=LINE, lw=0.5)
            text(s, Emu(xx + Inches(0.04)), Emu(yy + Inches(0.04)), Emu(cw - Inches(0.06)), Emu(h - Inches(0.06)),
                 str(val), size=sizes[ci], color=fg, bold=(ri == 0), anchor=MSO_ANCHOR.MIDDLE)
            xx = Emu(xx + cw)
        yy = Emu(yy + h)
    return yy


def add_pic(s, path, x, y, w, h):
    if os.path.exists(path):
        s.shapes.add_picture(path, x, y, width=w, height=h)
    else:
        rect(s, x, y, w, h, fill=ACC, line=LINE)


def build():
    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
    TOTAL = 24

    # 1 cover
    s = slide(prs)
    text(s, ML, Inches(1.35), Inches(12), Inches(0.3), "上海创智汇 × 同浦汇｜活动专项交付（双会资源融合）", size=14, color=BLUE, bold=True)
    text(s, ML, Inches(1.8), Inches(12), Inches(0.55), "创智汇 · 30场定义核与资源导入", size=32, color=INK, bold=True)
    text(s, ML, Inches(2.45), Inches(12), Inches(0.35), "AI+数字内容无界共创港　｜　2026.08—2027.07", size=15, color=GOLD, bold=True)
    text(s, ML, Inches(2.9), Inches(12), Inches(0.35),
         "每场≤30人　·　负责人≤30%　·　我们带客 / 园区销售　·　只融合招商贴近资源，不堆全量名单", size=13, color=MUT)
    card(s, ML, Inches(3.5), Inches(6.0), Inches(3.15), "创作主题响应", [
        "3F ≈2850㎡：AI / OPC / 算力 / 具身入孵",
        "5F ≈3670㎡：数字内容 / IP / 展贸转化",
        "WAIC 线喂养 3F 能力场次",
        "ChinaJoy 线喂养 5F 内容场次",
        "具身 / AIGC 为双会叠加带",
    ], BLUE)
    card(s, Inches(7.2), Inches(3.5), Inches(5.55), Inches(3.15), "本册边界", [
        "仅活动专项；不做租金租赁对赌",
        "双会表只取可导入切片，不粘贴全库",
        "出海 / 领事另议另计价",
        "不承诺必带外资企业",
        "主体：云创基地；学术：复旦住房中心",
    ], GREEN)
    footer(s, 1, TOTAL)

    # 2 主体
    s = slide(prs); header(s, "01 · 主体与支持", "谁在做这件事")
    rows = [
        ["角色", "单位", "说明"],
        ["主体运营范围", "上海市云计算创新基地", "国家级孵化器，纳入创智汇活动与运营主体范围"],
        ["学术支持", "复旦大学住房政策研究中心", "学术背书、议题与人才/空间相关研究支持"],
        ["活动及运营载体支持", "杨浦区科技企业联合会", "活动生态、会员协同、传播支持"],
        ["活动及运营载体支持", "杨浦区科技企业服务中心", "政策服务、申报辅导、入驻手续协同"],
        ["活动执行与带客", "同浦汇", "策划、邀约、执行、建群、带客回访"],
        ["销售成交", "园区", "房源销售、价格确认、合同交割"],
    ]
    table(s, Inches(1.35), rows, [0.22, 0.32, 0.46], sizes=[12, 14, 13], rh=0.55)
    footer(s, 2, TOTAL)

    # 3 口径
    s = slide(prs); header(s, "01 · 关键口径", "人数 · 负责人 · 客群 · 外企原则")
    card(s, ML, Inches(1.35), Inches(4.0), Inches(5.4), "人数与负责人", [
        "每场到场人次控制在 30 人以下",
        "全年触达目标约 600 人次+，以可执行、可核验为准",
        "负责人占比≤30%",
        "邀约以园区产业主题客群为主（弱化关联度表述）",
    ], BLUE)
    card(s, Inches(4.75), Inches(1.35), Inches(4.0), Inches(5.4), "外企与领事", [
        "原则：不承诺一定会带外资企业来",
        "出海活动：另议、价格另议，不在本册费用内",
        "领事到访接待：另计价",
        "领事馆挂牌及挂牌活动：另计价",
        "挂牌与领事活动性质不同，分开计价",
    ], GOLD)
    card(s, Inches(8.95), Inches(1.35), Inches(3.8), Inches(5.4), "分工一句话", [
        "我们负责带客",
        "园区负责销售部分",
        "活动结束建群",
        "园区可在群内促成成交",
    ], GREEN)
    footer(s, 3, TOTAL)

    # 4 主题响应 + 层次
    s = slide(prs); header(s, "02 · 内容层次", "响应「AI+数字内容无界共创港」的四层叠加",
                           "不是把双会全量搬进园区，而是按招商转化裁切")
    rows = [
        ["层次", "作用", "主要来源", "对应场次类型"],
        ["L0 点火", "渠道/金融/投研/政府点火", "园区自建+载体", "P1–P4 项目推广日"],
        ["L1 双会承接", "大会余热→园区看场", "WAIC成果切片", "E1 开放交流日"],
        ["L2 楼层主线", "3F能力 / 5F内容分轨获客", "WAIC赛道 + CJ优先策略", "A/B 与 C/D 常规场"],
        ["L3 双会叠加带", "具身+AIGC跨楼层联动", "WAIC具身×CJ Vision Future", "A5 等叠加体验"],
        ["L4 收口/预热", "签约收口与下届双会预热", "年内导入精选", "F4 / D5"],
    ]
    table(s, Inches(1.3), rows, [0.14, 0.28, 0.28, 0.3], sizes=[12, 12, 12, 12], rh=0.55)
    text(s, ML, Inches(5.5), CW, Inches(0.9),
         "原则：双会总表是「资源库」，30场是「定义核」。每场只导入可邀约、可看场、可建群的切片；不出全量展商名录。",
         size=13, color=MUT)
    footer(s, 4, TOTAL)

    # 5 双会资源摘要（招商贴近）
    s = slide(prs); header(s, "02 · 双会资源导入（摘要）", "只取与招商贴近的可调用规模与优先策略",
                           "数据来源：WAIC2026全量资源整合总表 · ChinaJoy2026招商引资总表六版融合")
    card(s, ML, Inches(1.3), Inches(5.9), Inches(5.45), "WAIC → 喂养 3F AI/OPC", [
        WAIC_POOL["scale"],
        "主赛道：算力芯片 / 机器人具身 / 大模型Agent / AIGC / 行业AI",
        "活动项目池缩尺：ACT-001闭门 · ACT-002具身 · ACT-003算力 · ACT-006行业AI",
        "用法：A/B线索抽样邀约；论坛议题变≤30人沙龙",
        "不做：963展商全量粘贴进PPT",
    ], BLUE)
    card(s, Inches(7.2), Inches(1.3), Inches(5.55), Inches(5.45), "ChinaJoy → 喂养 5F 内容/IP", [
        f"主题「{CJ_POOL['theme']}」· {CJ_POOL['scale']}",
        "最高优先：中小游戏工作室池（Express/Next Play/AGS）",
        "高优先：具身 Vision Future · 潮玩谷子供 · 创作者经济",
        CJ_POOL["fit"],
        "不做：422/966全量名单进PPT；大厂总部搬迁不作目标",
    ], GOLD)
    footer(s, 5, TOTAL)

    # 6-7 总表
    for subset, title, pg in ((EVENTS[:15], "30场概况总表（上）", 6), (EVENTS[15:], "30场概况总表（下）", 7)):
        s = slide(prs); header(s, "03 · 活动概况", title, "主题 · 时间 · 层次 · 资源源 · 价值落点")
        rows = [["编号", "主题/活动", "时间", "层次", "资源源", "价值落点"]]
        for e in subset:
            rows.append([e["code"], e["name"][:14], e["month"].replace("2026年", "").replace("2027年", "次"),
                         e["layer"][:8], e["src"][:12], e["lease"][:12]])
        table(s, Inches(1.25), rows, [0.08, 0.28, 0.14, 0.14, 0.18, 0.18], sizes=[9, 10, 9, 9, 9, 9], rh=0.32)
        footer(s, pg, TOTAL)

    # 8-9 定义核映射
    for subset, title, pg in ((EVENTS[:15], "30场定义核与资源导入（上）", 8), (EVENTS[15:], "30场定义核与资源导入（下）", 9)):
        s = slide(prs); header(s, "03 · 定义核", title, "每场一句话定义 + 导入什么资源（招商向）")
        rows = [["编号", "定义核（一句话）", "资源导入"]]
        for e in subset:
            rows.append([e["code"], e["define"][:28], e["import_res"][:36]])
        table(s, Inches(1.25), rows, [0.08, 0.42, 0.5], sizes=[9, 10, 10], rh=0.32)
        footer(s, pg, TOTAL)

    # 10 价值
    s = slide(prs); header(s, "04 · 活动价值", "活动带来什么")
    card(s, ML, Inches(1.35), Inches(5.9), Inches(5.4), "园区侧", [
        "每月有主题不断档，空间有热度",
        "8月项目推广日打通中介/银行/投研/政府",
        "3F/5F 分轨：AI入孵 + 内容展贸",
        "双会资源切片持续喂养带客池",
        "云创基地国家级孵化器背书可感知",
    ], BLUE)
    card(s, Inches(7.2), Inches(1.35), Inches(5.55), Inches(5.4), "转化侧", [
        "每场≤30人，精而准",
        "固定看场环节，便于带客",
        "活动结束建群，持续触达",
        "同浦汇带客，园区做销售促成",
        "不做租金租赁对赌与必要性要求",
    ], GREEN)
    footer(s, 10, TOTAL)

    # 11 交付
    s = slide(prs); header(s, "04 · 具体交付", "我们交给园区的东西")
    rows = [
        ["周期", "交付物", "说明"],
        ["签约后2周", "年度活动方案细化版 + 排期日历 + 资源导入台账模板", "活动专项"],
        ["每场结束后7天", "执行卡、议程、嘉宾确认、报名页、签到表、意向表、动线图、照片/通稿", "同浦汇归档"],
        ["每场结束后7天", "回访摘要（知会版）+ 建群清单 + 本场资源来源标注", "便于园区销售跟进"],
        ["每月", "活动数据月报 + 下月排期确认单", "联席会前"],
        ["每季/年末", "季度复盘、全年台账、双会资源消耗与补库建议", "资产沉淀"],
    ]
    table(s, Inches(1.35), rows, [0.18, 0.52, 0.3], sizes=[12, 12, 12], rh=0.58)
    text(s, ML, Inches(5.9), CW, Inches(0.7),
         "回访与带客由同浦汇负责；园区聚焦销售促成。交园区的回访材料为摘要知会，方便销售跟进，不额外增加园区事务负担。",
         size=12, color=MUT)
    footer(s, 11, TOTAL)

    # 12 分工
    s = slide(prs); header(s, "04 · 执行分工", "我们带客，园区负责销售")
    rows = [
        ["事项", "同浦汇", "园区", "载体支持/学术"],
        ["策划与议程", "全案策划、执行卡、资源切片", "方案审定、档期确认", "联合会/服中心协同"],
        ["邀约与带客", "双会线索抽样邀约、带客、建群", "共享在谈名单", "会员/学者推荐"],
        ["场地保障", "提前提报需求、物料", "报备、导视、物业安保、样板间", "—"],
        ["现场执行", "主持控场、看场带队", "现场支持、开放看场", "政策宣讲位（如需）"],
        ["销售促成", "提供意向名单与群", "负责销售部分、报价成交、合同", "入驻手续（服中心）"],
        ["回访跟踪", "主责带客回访与台账", "基于名单/群做销售跟进", "—"],
        ["复盘", "单场复盘+月报+资源补库", "确认下月活动安排", "—"],
    ]
    table(s, Inches(1.3), rows, [0.14, 0.3, 0.3, 0.26], sizes=[11, 11, 11, 11], rh=0.5)
    footer(s, 12, TOTAL)

    # 13 协同
    s = slide(prs); header(s, "04 · 协同着眼点", "做完活动、建完群 → 园区销售促成")
    card(s, ML, Inches(1.35), Inches(5.9), Inches(5.4), "标准动作", [
        "D日：同浦汇签到分层并建群",
        "T+1：同浦汇发看场预约与资料入群",
        "T+1~7：同浦汇持续带客回访",
        "园区：在群内发布房源/政策/条件，负责销售促成",
        "双方月度核对意向与成交进展",
    ], GREEN)
    card(s, Inches(7.2), Inches(1.35), Inches(5.55), Inches(5.4), "责任边界", [
        "我们负责带客",
        "园区负责销售部分",
        "回访摘要交园区便于销售跟进",
        "不承诺必带外资企业",
        "出海/领事相关另议另计价",
    ], BLUE)
    footer(s, 13, TOTAL)

    # 14 收费付款
    s = slide(prs); header(s, "05 · 收费与付款", "活动费用边界 · 付款节点（唯一方案）")
    rows = [
        ["项目", "口径", "说明"],
        ["活动策划执行服务费", "年度打包（金额线下确认）", "含常规30场策划执行、物料、主持、建群与带客回访"],
        ["租金参考", "3.3 元/㎡/天", "高于周边均价，故不做租赁对赌与租赁必要性要求"],
        ["建议免租期", "建议 1–3 个月（面议）", "降低决策门槛"],
        ["物业参考", "13.8 元/㎡/月", "园区统一标准"],
        ["待租房源", "8 间优先看场", "带客当场可看"],
        ["出海活动", "另议 · 另计价", "不在本册活动服务费内"],
        ["领事到访接待", "另计价", "在现有三部分费用之外"],
        ["领事馆挂牌及挂牌活动", "另计价", "与领事到访分属不同性质，分开计价"],
    ]
    table(s, Inches(1.2), rows, [0.24, 0.28, 0.48], sizes=[11, 12, 11], rh=0.38)
    text(s, ML, Inches(5.35), CW, Inches(0.28), "付款节点（唯一建议方案）", size=13, color=BLUE, bold=True)
    rows2 = [
        ["节点", "比例", "对应"],
        ["签约后7日内", "50% 启动款", "启动策划、年度日历锁定"],
        ["Q3 / Q4 / 次年Q1 / 次年Q2 各季", "各10%（合计40%）", "与季度场次完成度、月报验收挂钩"],
        ["年终收官后", "10% 尾款", "年报、台账、影像资产移交"],
    ]
    table(s, Inches(5.65), rows2, [0.4, 0.28, 0.32], sizes=[11, 11, 11], rh=0.34)
    footer(s, 14, TOTAL)

    # 15 执行节奏
    s = slide(prs); header(s, "06 · 单场执行节奏", "同一套打法")
    rows = [
        ["节点", "动作", "责任", "交付"],
        ["T-14", "锁定嘉宾议程物料+本场资源切片", "同浦汇", "议程/嘉宾确认/报名"],
        ["T-7", "邀约确认+场地报备", "同浦汇；园区场地确认", "名单/场地单"],
        ["T-1", "彩排+看场动线", "同浦汇；园区样板间", "彩排记录"],
        ["D日", "签到→主区→看场→建群", "同浦汇执行；园区支持", "签到/意向/群"],
        ["T+1~7", "带客回访", "同浦汇主责", "回访摘要知会园区"],
        ["销售促成", "房源条件与成交", "园区负责销售部分", "合同/交割"],
    ]
    table(s, Inches(1.35), rows, [0.12, 0.3, 0.32, 0.26], sizes=[11, 12, 11, 11], rh=0.55)
    footer(s, 15, TOTAL)

    # 16-17 cases
    for i, c in enumerate(CASES):
        s = slide(prs); header(s, f"07 · 往期案例（{i+1}）", c["name"], c["meta"])
        for j, fn in enumerate(c["photos"][:3]):
            add_pic(s, os.path.join(PHOTO, fn), Emu(ML + j * Inches(4.1)), Inches(1.35), Inches(3.95), Inches(2.3))
        rect(s, ML, Inches(3.85), Inches(7.5), Inches(2.95), fill=BG, line=LINE)
        text(s, Emu(ML + Inches(0.2)), Inches(3.95), Inches(7.1), Inches(0.3), "当天议程", size=13, color=BLUE, bold=True)
        bullets(s, Emu(ML + Inches(0.2)), Inches(4.35), Inches(7.1), Inches(2.2), c["agenda"], size=13)
        rect(s, Inches(8.3), Inches(3.85), Inches(4.5), Inches(2.95), fill=ACC, line=LINE)
        text(s, Inches(8.45), Inches(3.95), Inches(4.2), Inches(0.3), "可见机构/企业", size=13, color=BLUE, bold=True)
        text(s, Inches(8.45), Inches(4.35), Inches(4.2), Inches(2.2),
             f"{c['orgs']}\n\n企业侧：{c['companies']}\n\n映射：{c['map']}", size=12, color=INK)
        footer(s, 16 + i, TOTAL)

    # 18-20 逐场
    chunks = [(EVENTS[:10], 18), (EVENTS[10:20], 19), (EVENTS[20:], 20)]
    for subset, pg in chunks:
        s = slide(prs); header(s, "08 · 逐场执行要素", f"{subset[0]['code']}–{subset[-1]['code']}｜行程/嘉宾/交付")
        rows = [["编号", "活动", "行程", "嘉宾", "交付"]]
        for e in subset:
            rows.append([e["code"], e["name"][:14], e["agenda"][:24], e["guests"][:14], "标准包+建群" if "另" not in e["deliver"] else "另案"])
        table(s, Inches(1.3), rows, [0.07, 0.22, 0.3, 0.26, 0.15], sizes=[9, 10, 10, 10, 10], rh=0.42)
        footer(s, pg, TOTAL)

    # 21 双会衔接（完善）
    s = slide(prs); header(s, "09 · 双会衔接落地", "活动同步衔接 WAIC 与 ChinaJoy（招商向）")
    card(s, ML, Inches(1.3), Inches(5.9), Inches(5.45), "WAIC 线 → 3F", [
        "E1 开放交流日承接大会余热",
        "Agent / 多模态 / 算力 / 具身 / 治理常态化",
        "ACT项目池缩尺为≤30人场次",
        "年底 Demo Day 收口精选意向",
        "D5 预热 WAIC 2027",
    ], BLUE)
    card(s, Inches(7.2), Inches(1.3), Inches(5.55), Inches(5.45), "ChinaJoy 线 → 5F", [
        "数字娱乐与IP授权对接",
        "中小工作室 / 创作者首发",
        "潮玩供应链 / 市集转正",
        "IP联名与沉浸式联展",
        "主题响应「与AI同游」×园区共创港",
    ], GOLD)
    footer(s, 21, TOTAL)

    # 22 另计价
    s = slide(prs); header(s, "09 · 另议另计价事项", "不在本册活动服务费内")
    rows = [
        ["事项", "计价方式", "性质说明"],
        ["出海专题活动", "另议 · 另计价", "不纳入本册年度活动包"],
        ["领事到访接待", "另计价", "外事接待类"],
        ["领事馆挂牌及挂牌活动", "另计价", "挂牌类；与领事到访分属不同性质"],
        ["外资企业到场", "不承诺", "原则：不承诺一定会带外资企业来"],
    ]
    table(s, Inches(1.4), rows, [0.28, 0.28, 0.44], sizes=[13, 13, 13], rh=0.7)
    footer(s, 22, TOTAL)

    # 23 结论
    s = slide(prs); header(s, "10 · 一页纸结论", "活动专项交付要点（含资源导入）")
    card(s, ML, Inches(1.35), Inches(4.0), Inches(5.4), "我们交付", [
        "30场定义核 + 资源导入映射",
        "L0–L4 内容层次与叠加",
        "8月项目推广日四场点火",
        "WAIC→3F / CJ→5F 分轨",
        "带客、建群、回访摘要",
        "往期案例（企业+议程+现场照）",
    ], BLUE)
    card(s, Inches(4.75), Inches(1.35), Inches(4.0), Inches(5.4), "请园区协同", [
        "场地物业与样板间",
        "共享在谈名单、确认档期",
        "负责销售部分与成交交割",
        "群内发布条件促成成交",
        "云创基地/联合会/服中心协同",
    ], GREEN)
    card(s, Inches(8.95), Inches(1.35), Inches(3.8), Inches(5.4), "商务口径", [
        "付款：50%+4×10%+10%",
        "租金不做对赌/必要性要求",
        "建议免租期1–3个月",
        "出海/领事另计价",
        "不承诺必带外企",
    ], GOLD)
    footer(s, 23, TOTAL)

    # 24 end
    s = slide(prs)
    text(s, ML, Inches(2.2), Inches(12), Inches(0.5), "谢谢审阅", size=36, color=INK, bold=True, align=PP_ALIGN.CENTER)
    text(s, ML, Inches(3.0), Inches(12), Inches(0.4), "创智汇 · 30场定义核与双会资源导入", size=18, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    text(s, ML, Inches(3.55), Inches(12), Inches(0.35), "AI+数字内容无界共创港 · 我们带客 · 园区销售", size=14, color=MUT, align=PP_ALIGN.CENTER)
    text(s, ML, Inches(4.2), Inches(12), Inches(0.9),
         "主体运营：上海市云计算创新基地（国家级孵化器）\n学术支持：复旦大学住房政策研究中心\n载体支持：杨浦区科技企业联合会 · 科技企业服务中心",
         size=13, color=SOFT, align=PP_ALIGN.CENTER)
    footer(s, 24, TOTAL)

    os.makedirs(ART, exist_ok=True)
    outs = [
        "创智汇30场活动专项交付-双会资源融合版.pptx",
        "chuangzhihui-30-events-waic-cj-fuse.pptx",
        "创智汇30场活动专项交付-修订版.pptx",
        "chuangzhihui-30-events-activity-only-rev.pptx",
        "创智汇30场活动具体方案.pptx",
        "chuangzhihui-30-events-plan.pptx",
    ]
    tmp = os.path.join(HERE, "_tmp_fuse.pptx")
    prs.save(tmp)
    for fn in outs:
        p = os.path.join(HERE, fn)
        shutil.copy2(tmp, p)
        shutil.copy2(p, os.path.join(ART, fn))
    os.remove(tmp)
    print("PPT slides:", len(prs.slides))

    # Excel
    wb = Workbook()
    thin = Border(left=Side(style="thin", color="D0D0D0"), right=Side(style="thin", color="D0D0D0"),
                  top=Side(style="thin", color="D0D0D0"), bottom=Side(style="thin", color="D0D0D0"))
    head = PatternFill("solid", fgColor="1E4D8C")
    alt = PatternFill("solid", fgColor="EEF2F8")

    def hdr(ws, headers):
        for j, h in enumerate(headers, 1):
            c = ws.cell(1, j, h); c.font = Font(name="微软雅黑", bold=True, color="FFFFFF", size=10)
            c.fill = head; c.alignment = Alignment(wrap_text=True, vertical="center"); c.border = thin

    def write(ws, data, widths, h=44):
        for i, row in enumerate(data, 2):
            for j, v in enumerate(row, 1):
                c = ws.cell(i, j, v); c.font = Font(name="微软雅黑", size=9)
                c.alignment = Alignment(wrap_text=True, vertical="center"); c.border = thin
                if i % 2 == 0: c.fill = alt
            ws.row_dimensions[i].height = h
        ws.row_dimensions[1].height = 28
        for i, w in enumerate(widths, 1):
            ws.column_dimensions[get_column_letter(i)].width = w

    ws = wb.active; ws.title = "01-活动概况"
    hdr(ws, ["编号", "活动名称", "时间", "形式", "人数", "负责人占比", "场地", "层次", "资源源", "价值落点", "行程", "嘉宾", "交付"])
    write(ws, [[e["code"], e["name"], e["month"], e["fmt"], e["people"], e["dm"], e["floor"], e["layer"], e["src"],
                e["lease"], e["agenda"], e["guests"], e["deliver"]] for e in EVENTS],
          [8, 30, 16, 12, 10, 10, 12, 12, 18, 22, 36, 28, 28], 50)

    ws_def = wb.create_sheet("02-定义核资源导入")
    hdr(ws_def, ["编号", "活动名称", "层次", "定义核", "资源源", "资源导入（招商向）", "价值落点"])
    write(ws_def, [[e["code"], e["name"], e["layer"], e["define"], e["src"], e["import_res"], e["lease"]] for e in EVENTS],
          [8, 32, 12, 40, 22, 48, 22], 52)

    ws2 = wb.create_sheet("03-双会资源摘要")
    hdr(ws2, ["会展", "指标/策略", "招商用法", "映射场次"])
    write(ws2, [
        ["WAIC", WAIC_POOL["scale"], "A/B线索抽样+议题缩尺，不粘贴全库", "E1/A2/A3/A4/A5/B2/B4/C4/F3/F4"],
        ["WAIC", "ACT-001 大模型闭门", "缩尺为≤30人 Agent/圆桌", "A2/A3/F3"],
        ["WAIC", "ACT-002 机器人体验供需", "体验日+看场", "A5/E1"],
        ["WAIC", "ACT-003 算力芯片峰会", "厂商联训+券务", "A4/B4"],
        ["WAIC", "ACT-006 行业AI案例", "路演/政策沙龙", "C4/B1"],
        ["ChinaJoy", CJ_POOL["scale"] + " · 主题" + CJ_POOL["theme"], "可触达263优先外呼抽样", "C1/E4/D2/C3/D3/C5/D4/D5"],
        ["ChinaJoy", "最高：中小游戏工作室", "Express/Next Play/AGS 孵化向", "C1/E4/D2"],
        ["ChinaJoy", "高：具身 Vision Future", "与WAIC具身叠加", "A5/E1"],
        ["ChinaJoy", "高：潮玩谷子供应链", "5F展位/市集转正", "C3/D3/C5"],
        ["ChinaJoy", "中高：创作者经济", "首发+渠道复用", "D2/D5"],
        ["原则", "中腰部灵活/不追大厂总部", "双会是资源库，30场是定义核", "全册"],
    ], [12, 36, 36, 28], 40)

    ws3 = wb.create_sheet("04-主体分工收费付款")
    hdr(ws3, ["模块", "内容"])
    write(ws3, [
        ["主体运营范围", PARTIES["ops"]],
        ["学术支持", PARTIES["academic"]],
        ["活动及运营载体支持", PARTIES["support"]],
        ["执行带客", PARTIES["exec"]],
        ["销售", "园区负责销售部分与合同交割"],
        ["人数", "每场≤30人；全年触达约600人次+，可执行可核验"],
        ["负责人占比", "≤30%"],
        ["客群", "以园区产业主题邀约为主（弱化关联度/核验表表述）"],
        ["外企原则", "不承诺一定会带外资企业来"],
        ["出海", "另议另计价，不在本册费用内"],
        ["领事到访", "另计价（三部分费用之外）"],
        ["领事挂牌/挂牌活动", "另计价；与到访分属不同性质"],
        ["租金", "3.3元/㎡/天高于周边均价，不做租赁对赌与必要性要求"],
        ["免租期建议", "1–3个月（面议）"],
        ["付款", "签约后7日50%；Q3/Q4/次年Q1/次年Q2各10%；年终10%"],
        ["双会衔接", "WAIC→3F能力线；ChinaJoy→5F内容线；具身/AIGC为叠加带"],
        ["资源原则", "双会总表只取招商贴近切片；不做全量名单堆砌"],
    ], [18, 72], 34)

    ws4 = wb.create_sheet("05-往期案例")
    hdr(ws4, ["案例", "时间地点", "机构", "企业", "议程", "照片", "映射"])
    write(ws4, [[c["name"], c["meta"], c["orgs"], c["companies"], " | ".join(c["agenda"]), "；".join(c["photos"]), c["map"]] for c in CASES],
          [32, 28, 28, 36, 40, 28, 24], 56)

    ws5 = wb.create_sheet("06-层次叠加说明")
    hdr(ws5, ["层次", "作用", "主要来源", "对应场次"])
    write(ws5, [
        ["L0 点火", "渠道/金融/投研/政府点火", "园区自建+载体", "P1–P4"],
        ["L1 双会承接", "大会余热→园区看场", "WAIC成果切片", "E1"],
        ["L2 楼层主线", "3F能力 / 5F内容分轨", "WAIC赛道 + CJ优先策略", "A/B 与 C/D"],
        ["L3 双会叠加带", "具身+AIGC跨楼层", "WAIC具身×CJ Vision Future", "A5 等"],
        ["L4 收口/预热", "签约收口与下届预热", "年内导入精选", "F4 / D5"],
        ["另计价", "出海/领事不进年度包", "另案", "F1 / L1"],
    ], [14, 28, 32, 24], 36)

    xouts = [
        "创智汇30场活动专项交付-双会资源融合版.xlsx",
        "chuangzhihui-30-events-waic-cj-fuse.xlsx",
        "创智汇30场-双会资源导入摘要.xlsx",
        "创智汇30场活动专项交付-修订版.xlsx",
        "chuangzhihui-30-events-activity-only-rev.xlsx",
        "创智汇30场活动方案排期表.xlsx",
        "chuangzhihui-30-events-schedule.xlsx",
    ]
    for fn in xouts:
        p = os.path.join(HERE, fn); wb.save(p); shutil.copy2(p, os.path.join(ART, fn))
    print("Excel OK")


if __name__ == "__main__":
    build()
