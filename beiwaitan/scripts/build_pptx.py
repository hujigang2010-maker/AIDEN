"""
北外滩 · 一滴水 · 人文思想之夜 · 品牌规划 PPT 生成脚本

设计原则：
- 严格按用户要求，不含王德峰个人相关内容
- 只体现"品牌 / 空间 / 主办联合体 / 内容形态 / 报名机制 / 系列规划"

视觉风格：
- 东方人文美学 · 极简 · 大留白
- 主色板：
    墨绿  #2C5F3E  主基调 · 竹林/山水
    深梅  #6B2C3E  强调 · 文人色
    金    #C9A24A  重点 · 金笺
    米白  #F7F1E6  底色 · 素笺
    炭黑  #1B1F2A  正文
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pathlib import Path

MOSS = RGBColor(0x2C, 0x5F, 0x3E)     # 墨绿
PLUM = RGBColor(0x6B, 0x2C, 0x3E)     # 深梅
GOLD = RGBColor(0xC9, 0xA2, 0x4A)     # 金
IVORY = RGBColor(0xF7, 0xF1, 0xE6)    # 米白
INK = RGBColor(0x1B, 0x1F, 0x2A)       # 炭黑
GREY = RGBColor(0x6B, 0x73, 0x80)      # 灰
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD8, 0xD0, 0xC0)      # 素笺边线
DEEP = RGBColor(0x1A, 0x3E, 0x27)      # 深墨绿

CN_FONT = "WenQuanYi Micro Hei"
EN_FONT = "Georgia"  # 衬线，配合人文气质

SLIDES = []


def set_run(run, text, *, size=14, bold=False, color=INK, italic=False,
            font_cn=CN_FONT, font_en=EN_FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_en
    rPr = run._r.get_or_add_rPr()
    for tag in ("ea", "cs"):
        existing = rPr.find(qn(f"a:{tag}"))
        if existing is not None:
            rPr.remove(existing)
        ea = rPr.makeelement(qn(f"a:{tag}"), {"typeface": font_cn})
        rPr.append(ea)


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, fill=None, line=None,
             italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if fill is not None:
        box.fill.solid(); box.fill.fore_color.rgb = fill
    else:
        box.fill.background()
    if line is not None:
        box.line.color.rgb = line; box.line.width = Pt(0.5)
    else:
        box.line.fill.background()
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        set_run(run, ln, size=size, bold=bold, color=color, italic=italic)
    return box


def add_rect(slide, x, y, w, h, *, fill=MOSS, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(0.5)
    shp.shadow.inherit = False
    shp.text_frame.text = ""
    return shp


def add_round(slide, x, y, w, h, text="", *, fill=PLUM, color=IVORY,
              size=12, bold=True, align=PP_ALIGN.CENTER):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = align
    set_run(p.add_run(), text, size=size, bold=bold, color=color)
    return shp


def add_oval(slide, x, y, w, h, *, fill=None, line_color=None, line_w=1.5):
    """圆——一滴水的涟漪母题"""
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_table(slide, x, y, w, h, header, rows, *, header_fill=MOSS,
              header_color=IVORY, zebra=(WHITE, IVORY), header_size=11,
              body_size=10, col_widths=None, body_align=PP_ALIGN.LEFT):
    cols = len(header); n_rows = len(rows) + 1
    ts = slide.shapes.add_table(n_rows, cols, x, y, w, h)
    table = ts.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw
    for j, htxt in enumerate(header):
        cell = table.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Pt(4); tf.margin_right = Pt(4)
        tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
        tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        set_run(p.add_run(), htxt, size=header_size, bold=True, color=header_color)
    for i, row in enumerate(rows, start=1):
        bg = zebra[(i - 1) % 2]
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Pt(4); tf.margin_right = Pt(4)
            tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
            tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = body_align
            set_run(p.add_run(), str(val), size=body_size, color=INK)
    return ts


def add_chrome(slide, prs, *, page_no, phase_label="", page_title="", subtitle=""):
    """页眉页脚 · 米白底 + 墨绿边"""
    sw, sh = prs.slide_width, prs.slide_height
    # 顶部三层：墨绿主条 + 金细条 + 米白留白
    add_rect(slide, 0, 0, sw, Emu(400000), fill=MOSS)
    add_rect(slide, 0, Emu(400000), sw, Emu(20000), fill=GOLD)
    if phase_label:
        add_round(slide, Inches(0.5), Inches(0.16), Inches(2.4), Inches(0.38),
                  phase_label, fill=GOLD, color=MOSS, size=11, bold=True)
    if page_title:
        add_text(slide, Inches(3.1), Inches(0.10), Inches(9.5), Inches(0.55),
                 page_title, size=22, bold=True, color=IVORY,
                 anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(3.1), Inches(0.55), Inches(9.5), Inches(0.28),
                 subtitle, size=11, color=IVORY, anchor=MSO_ANCHOR.MIDDLE)
    # 底部：金细条
    add_rect(slide, 0, sh - Emu(60000), sw, Emu(60000), fill=GOLD)
    add_text(slide, Inches(0.5), Inches(7.05), Inches(9.0), Inches(0.30),
             "北外滩 · 一滴水 · 人文思想之夜 · 品牌规划 v1.0",
             size=9, color=GREY)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(2.0), Inches(0.30),
             f"{page_no} / 0", size=9, color=GREY, align=PP_ALIGN.RIGHT)


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    SLIDES.append(s)
    return s


# ============================ 主流程 ============================
def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    sw, sh = prs.slide_width, prs.slide_height

    # ================= 1. 封面 =================
    s = new_slide(prs)
    # 底色米白
    add_rect(s, 0, 0, sw, sh, fill=IVORY)
    # 顶部墨绿细条
    add_rect(s, 0, 0, sw, Emu(80000), fill=MOSS)
    # 底部金细条
    add_rect(s, 0, sh - Emu(80000), sw, Emu(80000), fill=GOLD)

    # 一滴水涟漪母题（右侧多重同心圆）
    ripple_cx = Inches(11.0); ripple_cy = Inches(3.75)
    for r in [2.8, 2.2, 1.6, 1.0, 0.5]:
        rad = Inches(r)
        add_oval(s, ripple_cx - rad, ripple_cy - rad, rad * 2, rad * 2,
                 line_color=MOSS if r != 0.5 else PLUM, line_w=0.8 if r > 0.7 else 2)

    # 大标题
    add_text(s, Inches(0.8), Inches(1.0), Inches(6.5), Inches(0.5),
             "BEIWAITAN · Humanities Nights", size=14, bold=True,
             color=MOSS, italic=True)
    add_text(s, Inches(0.8), Inches(1.8), Inches(9), Inches(1.5),
             "北外滩 · 一滴水", size=52, bold=True, color=MOSS)
    add_text(s, Inches(0.8), Inches(3.0), Inches(9), Inches(1.5),
             "人文思想之夜", size=52, bold=True, color=PLUM)

    # 金色分隔线
    add_rect(s, Inches(0.8), Inches(4.5), Inches(2.5), Emu(30000), fill=GOLD)

    # 品牌一句话
    add_text(s, Inches(0.8), Inches(4.8), Inches(9), Inches(0.5),
             "在时代变化中寻找生命的意义与智慧",
             size=20, color=INK, italic=True)

    # 副信息
    add_text(s, Inches(0.8), Inches(5.7), Inches(9), Inches(0.4),
             "融合 · 哲学思想 / 城市发展 / 企业家精神 / 人文交流",
             size=13, color=GREY)

    add_text(s, Inches(0.8), Inches(6.15), Inches(9), Inches(0.4),
             "首场 · 2026.10.31（周六）· 上海北外滩 · 一滴水",
             size=13, color=MOSS, bold=True)

    add_text(s, Inches(0.8), Inches(6.65), Inches(9), Inches(0.35),
             "品牌规划 v1.0 · 汇报对象：主办联合体三方",
             size=10, color=GREY)

    # ================= 2. 议程 =================
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="议程",
               page_title="汇报议程",
               subtitle="从「一场活动」到「长期人文品牌」的规划路径")
    items = [
        ("01", "核心一页 · 品牌定位与三年愿景", "One-Pager"),
        ("02", "落地空间 · 北外滩 · 一滴水", "Venue"),
        ("03", "主办联合体 · 三方协同", "Partners"),
        ("04", "内容框架 · 下午 + 晚间双段式", "Format"),
        ("05", "邀请对象 · 五类嘉宾矩阵", "Guests"),
        ("06", "系列化路径 · 三年品牌演进", "Roadmap"),
        ("07", "预算 / 组织 / 传播 / 风险", "Operations"),
    ]
    y0 = Inches(1.3)
    for i, (no, title, sub) in enumerate(items):
        y = y0 + Inches(0.75) * i
        add_round(s, Inches(0.8), y, Inches(0.7), Inches(0.55), no,
                  fill=MOSS, color=IVORY, size=18, bold=True)
        add_text(s, Inches(1.7), y, Inches(6.5), Inches(0.55), title,
                 size=17, bold=True, color=MOSS, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(8.2), y, Inches(5), Inches(0.55), sub,
                 size=13, color=GREY, italic=True, anchor=MSO_ANCHOR.MIDDLE)
        # 米白分割线
        add_rect(s, Inches(0.8), y + Inches(0.62), Inches(11.7), Emu(15000), fill=LINE)

    # ================= 3. 一页看懂 =================
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="01 · 核心一页",
               page_title="一页看懂 · 品牌定位与三年愿景",
               subtitle="融合哲学 · 城市 · 企业家 · 人文 的高品质思想交流平台")

    # 上部 · 品牌一句话
    add_rect(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(1.2), fill=MOSS)
    add_text(s, Inches(0.5), Inches(1.20), Inches(12.3), Inches(0.4),
             "品牌一句话", size=11, color=GOLD, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.6),
             "在时代变化中寻找生命的意义与智慧",
             size=26, color=IVORY, bold=True, italic=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 中部 · 四大关键词
    keywords = [
        ("哲学思想", "启迪智慧", MOSS),
        ("城市发展", "关照当下", PLUM),
        ("企业家精神", "驱动未来", GOLD),
        ("人文交流", "共话生命", DEEP),
    ]
    kw_y = Inches(2.55); kw_w = Inches(2.95); kw_h = Inches(1.5); kw_gap = Inches(0.15)
    for i, (w, sub, c) in enumerate(keywords):
        x = Inches(0.5) + (kw_w + kw_gap) * i
        add_rect(s, x, kw_y, kw_w, kw_h, fill=c)
        add_text(s, x, kw_y + Inches(0.20), kw_w, Inches(0.6), w,
                 size=22, bold=True,
                 color=GOLD if c == MOSS else IVORY,
                 align=PP_ALIGN.CENTER)
        add_text(s, x, kw_y + Inches(0.85), kw_w, Inches(0.5), sub,
                 size=14, color=IVORY, italic=True,
                 align=PP_ALIGN.CENTER)

    # 下部 · 三年愿景
    add_text(s, Inches(0.5), Inches(4.30), Inches(12.3), Inches(0.4),
             "三年演进 · 从首场标杆到品牌矩阵",
             size=14, bold=True, color=MOSS)
    add_rect(s, Inches(0.5), Inches(4.72), Inches(12.3), Emu(15000), fill=GOLD)

    vision = [
        ("Y1 · 2026", "试运行", "首场（10.31）+ 第二场\n验证形态 · 树立标杆", MOSS),
        ("Y2 · 2027", "品牌立起来", "全年 4 场\n联合出品 · 长期社群", PLUM),
        ("Y3 · 2028", "品牌矩阵", "季度大会 + 月度沙龙\n形成「北外滩思想系列」", GOLD),
    ]
    vy = Inches(4.90); vw = Inches(4.0); vh = Inches(1.85)
    for i, (y_label, tag, desc, c) in enumerate(vision):
        x = Inches(0.5) + (vw + Inches(0.15)) * i
        add_rect(s, x, vy, vw, vh, fill=IVORY, line=LINE)
        add_rect(s, x, vy, vw, Inches(0.5), fill=c)
        add_text(s, x, vy + Inches(0.05), vw, Inches(0.4), y_label,
                 size=15, bold=True,
                 color=MOSS if c == GOLD else IVORY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, vy + Inches(0.55), vw, Inches(0.35), tag,
                 size=13, bold=True, color=c, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.15), vy + Inches(0.95), vw - Inches(0.3),
                 vh - Inches(1.0), desc, size=11, color=INK, align=PP_ALIGN.CENTER)

    # 底部宣示
    add_text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.35),
             "打造具有辨识度的城市人文思想 IP · 让思想在北外滩生根",
             size=11, italic=True, color=GREY, align=PP_ALIGN.CENTER)

    # ================= 4. 品牌命名与主张 =================
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="01 · 品牌",
               page_title="品牌命名与主张",
               subtitle="品牌命名的哲学 · 一滴水见汪洋")

    # 中间大字
    add_text(s, Inches(0.8), Inches(1.30), Inches(6), Inches(0.4),
             "系列品牌名", size=12, color=GREY, italic=True)
    add_text(s, Inches(0.8), Inches(1.75), Inches(6), Inches(0.7),
             "北外滩 · 一滴水", size=32, bold=True, color=MOSS)
    add_text(s, Inches(0.8), Inches(2.55), Inches(6), Inches(0.7),
             "人文思想之夜", size=32, bold=True, color=PLUM)

    # 命名释义
    add_rect(s, Inches(0.8), Inches(3.60), Inches(6), Emu(20000), fill=GOLD)
    add_text(s, Inches(0.8), Inches(3.75), Inches(6), Inches(0.4),
             "命名释义", size=13, bold=True, color=MOSS)
    add_text(s, Inches(0.8), Inches(4.15), Inches(6), Inches(2.6),
             "北外滩 —— 上海面向世界的门户，"
             "承载着近代以来的思想启蒙与文化交融\n\n"
             "一滴水 —— 见微知著，见汪洋大海。"
             "以「小空间」承载「大思考」，"
             "以「一场活动」积淀「长期品牌」\n\n"
             "人文思想之夜 —— 不同于产业峰会的白日喧嚣，"
             "选择夜的深邃，让思想在从容中生根",
             size=13, color=INK)

    # 右侧 · 一滴水视觉母题
    ripple_cx = Inches(10.0); ripple_cy = Inches(4.0)
    for r in [2.5, 1.9, 1.3, 0.7, 0.25]:
        rad = Inches(r)
        add_oval(s, ripple_cx - rad, ripple_cy - rad, rad * 2, rad * 2,
                 line_color=MOSS if r != 0.25 else PLUM,
                 line_w=1.5 if r > 0.5 else 3, fill=None if r > 0.3 else PLUM)
    add_text(s, Inches(7.5), Inches(6.6), Inches(5), Inches(0.4),
             "一滴水见汪洋 · 涟漪母题",
             size=11, italic=True, color=GREY, align=PP_ALIGN.CENTER)

    # ================= 5. 落地空间 =================
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="02 · 落地空间",
               page_title="北外滩 · 一滴水",
               subtitle="上海市虹口区北外滩东大名路 501 号")

    # 左：地址与空间描述
    add_text(s, Inches(0.5), Inches(1.15), Inches(6), Inches(0.4),
             "空间简介", size=14, bold=True, color=MOSS)
    add_rect(s, Inches(0.5), Inches(1.55), Inches(6), Inches(4.8),
             fill=IVORY, line=LINE)
    add_text(s, Inches(0.7), Inches(1.70), Inches(5.6), Inches(4.6),
             "「一滴水」位于上海北外滩核心地段，"
             "毗邻黄浦江，与外滩万国建筑群隔江相望，"
             "俯瞰上海城市历史与未来。\n\n"
             "空间设计以「水」为灵感，融合现代极简"
             "与东方美学，适合承办思想类沙龙、"
             "私享晚宴与文化对话。\n\n"
             "地址：上海市虹口区北外滩东大名路 501 号\n\n"
             "交通：\n"
             "  • 地铁 12 号线国际客运中心站\n"
             "  • 距外滩步行 15 分钟\n"
             "  • 距南京东路 10 分钟车程\n\n"
             "适用场景：\n"
             "  • 半日思想分享（可容 100–200 人）\n"
             "  • 私享晚宴（40–60 人圆桌）",
             size=11, color=INK)

    # 右：空间气质关键词卡片
    add_text(s, Inches(7.0), Inches(1.15), Inches(6), Inches(0.4),
             "为什么选一滴水", size=14, bold=True, color=MOSS)
    reasons = [
        ("城市地标", "北外滩 · 上海面向世界的新窗口", MOSS),
        ("水的意境", "极简东方美学 · 契合「一滴水」品牌", PLUM),
        ("场景灵活", "半日分享 + 晚间私享 · 一站式承接", GOLD),
        ("交通便捷", "地铁 + 轮渡 + 步行外滩", DEEP),
        ("传播价值", "背景即品牌 · 摄影/直播天然出片", MOSS),
    ]
    for i, (t, d, c) in enumerate(reasons):
        y = Inches(1.55) + Inches(0.95) * i
        add_rect(s, Inches(7.0), y, Inches(0.20), Inches(0.85), fill=c)
        add_rect(s, Inches(7.25), y, Inches(5.65), Inches(0.85),
                 fill=IVORY, line=LINE)
        add_text(s, Inches(7.4), y + Inches(0.08), Inches(2.0), Inches(0.35),
                 t, size=13, bold=True, color=MOSS, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(7.4), y + Inches(0.42), Inches(5.4), Inches(0.4),
                 d, size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # ================= 6. 主办联合体 =================
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="03 · 主办",
               page_title="主办联合体 · 三方协同",
               subtitle="学术公信力 + 品牌运营 + 产业与政府接口 三足鼎立")

    partners = [
        ("复旦大学\n住房政策研究中心", "学术背书",
         ["顶级学者资源", "内容深度把关", "学术公信力", "研究成果输出"],
         MOSS),
        ("首乾书院", "品牌运营",
         ["品牌主理与设计", "内容策划与执行", "现场执行团队",
          "长期品牌运营"], PLUM),
        ("上海市杨浦区\n科技企业联合会", "产业与政府接口",
         ["企业家资源池", "政府关系协同", "传播渠道联动", "场地/交通协调"],
         GOLD),
    ]
    cw = Inches(4.10); cy = Inches(1.20); ch = Inches(5.5); gx = Inches(0.10)
    for i, (name, role, contribs, c) in enumerate(partners):
        x = Inches(0.5) + (cw + gx) * i
        add_rect(s, x, cy, cw, ch, fill=IVORY, line=LINE)
        add_rect(s, x, cy, cw, Inches(1.5), fill=c)
        add_text(s, x, cy + Inches(0.20), cw, Inches(0.9), name,
                 size=17, bold=True,
                 color=MOSS if c == GOLD else IVORY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, cy + Inches(1.05), cw, Inches(0.35), role,
                 size=13, italic=True,
                 color=MOSS if c == GOLD else IVORY, align=PP_ALIGN.CENTER)

        add_text(s, x + Inches(0.15), cy + Inches(1.70), cw - Inches(0.3),
                 Inches(0.4), "主要贡献", size=11, bold=True,
                 color=c, align=PP_ALIGN.LEFT)
        for j, cn in enumerate(contribs):
            y = cy + Inches(2.15) + Inches(0.60) * j
            add_round(s, x + Inches(0.25), y + Inches(0.10), Inches(0.30),
                      Inches(0.30), str(j + 1), fill=c,
                      color=MOSS if c == GOLD else IVORY, size=11, bold=True)
            add_text(s, x + Inches(0.65), y + Inches(0.03), cw - Inches(0.8),
                     Inches(0.45), cn, size=12, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # ================= 7. 活动内容框架 =================
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="04 · 内容",
               page_title="活动内容框架 · 下午 + 晚间双段式",
               subtitle="半日思想分享 → 晚间人文私享晚宴 · 标准可复制模式")

    # 左：下午议程
    add_text(s, Inches(0.5), Inches(1.15), Inches(6), Inches(0.4),
             "下午 · 思想分享（14:00–17:30）",
             size=15, bold=True, color=MOSS)
    add_rect(s, Inches(0.5), Inches(1.60), Inches(6), Inches(4.2),
             fill=IVORY, line=LINE)
    afternoon = [
        ("14:00–14:30", "嘉宾签到"),
        ("14:30–14:40", "开场致辞"),
        ("14:40–16:10", "主题分享（90 分钟）"),
        ("16:10–16:40", "互动交流"),
        ("16:40–17:30", "圆桌对话与交流"),
    ]
    for i, (t, d) in enumerate(afternoon):
        y = Inches(1.75) + Inches(0.72) * i
        add_rect(s, Inches(0.65), y, Inches(1.65), Inches(0.60), fill=MOSS)
        add_text(s, Inches(0.65), y, Inches(1.65), Inches(0.60), t,
                 size=11, bold=True, color=IVORY, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(2.40), y, Inches(3.95), Inches(0.60), d,
                 size=12, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # 右：晚间议程
    add_text(s, Inches(7.0), Inches(1.15), Inches(6), Inches(0.4),
             "晚间 · 人文私享晚宴（18:00–21:00）",
             size=15, bold=True, color=PLUM)
    add_rect(s, Inches(7.0), Inches(1.60), Inches(6), Inches(4.2),
             fill=IVORY, line=LINE)
    evening = [
        ("18:00–18:30", "嘉宾签到"),
        ("18:30–18:40", "开场致辞"),
        ("18:40–20:30", "人文私享晚宴（圆桌交流）"),
        ("20:30–21:00", "自由交流"),
    ]
    for i, (t, d) in enumerate(evening):
        y = Inches(1.75) + Inches(0.90) * i
        add_rect(s, Inches(7.15), y, Inches(1.65), Inches(0.75), fill=PLUM)
        add_text(s, Inches(7.15), y, Inches(1.65), Inches(0.75), t,
                 size=11, bold=True, color=IVORY, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(8.90), y, Inches(3.95), Inches(0.75), d,
                 size=12, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # 底部：三个亮点关键词
    add_rect(s, Inches(0.5), Inches(6.05), Inches(12.5), Inches(0.85), fill=MOSS)
    highlights = [("思想分享", "启迪智慧"), ("高端交流", "链接未来"),
                  ("私享晚宴", "共话发展")]
    for i, (t, sub) in enumerate(highlights):
        x = Inches(0.5) + Inches(4.17) * i
        add_text(s, x, Inches(6.10), Inches(4.15), Inches(0.4), t,
                 size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(6.50), Inches(4.15), Inches(0.4), sub,
                 size=12, color=IVORY, italic=True, align=PP_ALIGN.CENTER)

    # ================= 8. 邀请对象矩阵 =================
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="05 · 邀请",
               page_title="邀请对象矩阵 · 五类嘉宾画像",
               subtitle="80–120 人 · 高质量小规模 · 审核邀请制")

    # 左：五类嘉宾卡
    guests = [
        ("高校学者", "复旦 / 交大 / 同济 / 华东师大 教授、研究员",
         20, MOSS),
        ("企业家", "上海及长三角优秀企业家（一代 / 二代）",
         35, PLUM),
        ("科技创新人士", "智能网联 / AI / 生物医药 / 新材料 创业者与高管",
         20, GOLD),
        ("文化领域嘉宾", "文创 / 出版 / 艺术 / 媒体主理人",
         15, DEEP),
        ("城市发展关注者", "建筑师 / 规划师 / 政府智库",
         10, MOSS),
    ]
    for i, (t, d, pct, c) in enumerate(guests):
        y = Inches(1.20) + Inches(1.10) * i
        add_rect(s, Inches(0.5), y, Inches(0.22), Inches(1.0), fill=c)
        add_rect(s, Inches(0.75), y, Inches(7.5), Inches(1.0),
                 fill=IVORY, line=LINE)
        add_text(s, Inches(0.9), y + Inches(0.05), Inches(2.5), Inches(0.35),
                 t, size=14, bold=True, color=c, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(0.9), y + Inches(0.45), Inches(6.9), Inches(0.5),
                 d, size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        # 比例圆
        add_round(s, Inches(7.4), y + Inches(0.28), Inches(0.85),
                  Inches(0.45), f"{pct}%",
                  fill=c, color=MOSS if c == GOLD else IVORY,
                  size=14, bold=True)

    # 右：邀请策略
    add_text(s, Inches(8.6), Inches(1.15), Inches(4.4), Inches(0.4),
             "邀请策略", size=14, bold=True, color=MOSS)

    strategy = [
        ("总人数", "80–120 人（首场）", MOSS),
        ("邀请方式", "审核邀请制（非公开）", PLUM),
        ("筛选标准", "行业代表性 + 表达欲望 + 长期性", GOLD),
        ("性别比例", "建议 55/45 或 60/40", DEEP),
        ("年龄结构", "30–60 岁核心 + 20% 青年学者", MOSS),
        ("VIP 席位", "预留 8–12 位（一线嘉宾+主办)", PLUM),
    ]
    for i, (t, d, c) in enumerate(strategy):
        y = Inches(1.60) + Inches(0.75) * i
        add_rect(s, Inches(8.6), y, Inches(0.15), Inches(0.65), fill=c)
        add_text(s, Inches(8.80), y + Inches(0.05), Inches(1.7), Inches(0.55),
                 t, size=11, bold=True, color=MOSS, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(10.5), y + Inches(0.05), Inches(2.55), Inches(0.55),
                 d, size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # ================= 9. 报名与筛选机制 =================
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="05 · 报名",
               page_title="报名与筛选 · 预报名 + 审核邀请制",
               subtitle="两阶段 · 沉淀社群资产 · 长期运营")

    # 4 步流程
    add_text(s, Inches(0.5), Inches(1.15), Inches(12), Inches(0.4),
             "四步流程", size=14, bold=True, color=MOSS)
    add_rect(s, Inches(0.5), Inches(1.55), Inches(12.3), Emu(15000), fill=GOLD)

    steps = [
        ("① 预报名", "扫码 → 填基础信息\n（姓名/单位/职务/一句话自我介绍）", MOSS),
        ("② 筛选审核", "主办方联合审核\n参照 5 类嘉宾配比", PLUM),
        ("③ 正式邀请", "审核通过发正式邀请函\n+ 详细议程 + 交通指引", GOLD),
        ("④ 社群沉淀", "所有报名者进活动预备群\n形成长期社群资产", DEEP),
    ]
    sw2 = Inches(2.95); sy = Inches(1.75); sh2 = Inches(2.0); sgap = Inches(0.15)
    for i, (t, d, c) in enumerate(steps):
        x = Inches(0.5) + (sw2 + sgap) * i
        add_rect(s, x, sy, sw2, sh2, fill=IVORY, line=LINE)
        add_rect(s, x, sy, sw2, Inches(0.6), fill=c)
        add_text(s, x, sy + Inches(0.10), sw2, Inches(0.4), t,
                 size=16, bold=True,
                 color=MOSS if c == GOLD else IVORY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.15), sy + Inches(0.75),
                 sw2 - Inches(0.3), sh2 - Inches(0.85),
                 d, size=11, color=INK)

    # 下部 · 数据资产说明
    add_text(s, Inches(0.5), Inches(4.20), Inches(12), Inches(0.4),
             "社群资产逻辑 · 从报名到长期沉淀",
             size=14, bold=True, color=MOSS)
    add_rect(s, Inches(0.5), Inches(4.60), Inches(12.3), Inches(2.2),
             fill=IVORY, line=LINE)
    assets = [
        "首场活动预报名 → 200–400 人核心社群",
        "微信社群按邀请对象分类：学者群 / 企业家群 / 科技群 / 文化群 / 综合群",
        "每场活动前 30 天启动新一轮预报名，老学员优先通过",
        "3 年沉淀目标：核心社群 2,000+ 人，形成「北外滩思想社群」",
        "长期变现路径：付费会员制（Y3 起）· 书籍/课程/内训 分销",
    ]
    for i, a in enumerate(assets):
        y = Inches(4.75) + Inches(0.40) * i
        add_round(s, Inches(0.75), y + Inches(0.10), Inches(0.22), Inches(0.22),
                  "●", fill=MOSS, color=IVORY, size=8, bold=True)
        add_text(s, Inches(1.05), y, Inches(11.7), Inches(0.4),
                 a, size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # ================= 10. 系列化路径（三年） =================
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="06 · 系列化",
               page_title="三年系列化路径",
               subtitle="从首场标杆 → 品牌立起来 → 品牌矩阵")

    # 三年时间轴（横向）
    years = [
        ("Y1 · 2026", "试运行 · 树标杆",
         ["10 月首场活动（10.31）",
          "12 月第二场（半年检验）",
          "沉淀首批 200-400 人社群",
          "形成首套「下午 + 晚间」标准模板"], MOSS),
        ("Y2 · 2027", "品牌立起来",
         ["全年 4 场（3/5/7/10 月）",
          "联合出品：书籍 / 视频号 / 长期社群",
          "首场「付费入场」探索",
          "社群扩至 1,000+ 人"], PLUM),
        ("Y3 · 2028", "品牌矩阵",
         ["每季度 1 场大会 + 每月 1 场闭门沙龙",
          "形成「北外滩思想系列」 IP",
          "延伸城市/校园/企业专场",
          "社群 2,000+ / 付费会员制启动"], GOLD),
    ]
    yw = Inches(4.10); yy = Inches(1.20); yh = Inches(5.5); ygap = Inches(0.10)
    for i, (y_label, tag, bullets, c) in enumerate(years):
        x = Inches(0.5) + (yw + ygap) * i
        add_rect(s, x, yy, yw, yh, fill=IVORY, line=LINE)
        add_rect(s, x, yy, yw, Inches(1.05), fill=c)
        add_text(s, x, yy + Inches(0.10), yw, Inches(0.5), y_label,
                 size=20, bold=True,
                 color=MOSS if c == GOLD else IVORY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, yy + Inches(0.62), yw, Inches(0.4), tag,
                 size=14, bold=True,
                 color=MOSS if c == GOLD else IVORY,
                 align=PP_ALIGN.CENTER, italic=True)
        for j, b in enumerate(bullets):
            by = yy + Inches(1.35) + Inches(0.85) * j
            add_round(s, x + Inches(0.20), by + Inches(0.18), Inches(0.30),
                      Inches(0.30), str(j + 1), fill=c,
                      color=MOSS if c == GOLD else IVORY,
                      size=12, bold=True)
            add_text(s, x + Inches(0.55), by, yw - Inches(0.7), Inches(0.85),
                     b, size=11, color=INK)

    # ================= 11. 首场活动落地清单 =================
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="06 · 首场",
               page_title="首场活动 · 10.31 · 落地清单",
               subtitle="从筹备到复盘 · 8 个关键里程碑")

    milestones = [
        ("T-90 天", "8 月上旬", "立项 + 三方联合体协议 + 场地档期确认"),
        ("T-75 天", "8 月中旬", "内容策划完成 + 主题方向确定 + 嘉宾邀请函发出"),
        ("T-60 天", "9 月上旬", "主视觉 + 海报 + 预报名通道上线"),
        ("T-45 天", "9 月中旬", "预报名启动 · 目标 400+ 人 · 分群运营"),
        ("T-30 天", "10 月初", "筛选审核 · 正式邀请 · 席位分配"),
        ("T-14 天", "10 月中", "彩排 · 直播/录制方案 · 晚宴菜单确认"),
        ("T-3 天", "10.28-30", "现场布置 · 全员彩排 · 应急预案"),
        ("T 日", "10.31 · 周六", "正式举办 + 24h 内出媒体稿"),
    ]
    add_text(s, Inches(0.5), Inches(1.15), Inches(12), Inches(0.4),
             "8 个关键里程碑（倒推时间轴）",
             size=14, bold=True, color=MOSS)
    add_rect(s, Inches(0.5), Inches(1.55), Inches(12.3), Emu(15000), fill=GOLD)

    for i, (t, date, desc) in enumerate(milestones):
        y = Inches(1.70) + Inches(0.60) * i
        add_rect(s, Inches(0.5), y, Inches(1.5), Inches(0.50), fill=MOSS)
        add_text(s, Inches(0.5), y, Inches(1.5), Inches(0.50), t,
                 size=12, bold=True, color=IVORY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(2.05), y, Inches(1.7), Inches(0.50),
                 fill=GOLD)
        add_text(s, Inches(2.05), y, Inches(1.7), Inches(0.50), date,
                 size=11, bold=True, color=MOSS,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(3.80), y, Inches(9.0), Inches(0.50),
                 fill=IVORY, line=LINE)
        add_text(s, Inches(3.95), y, Inches(8.8), Inches(0.50), desc,
                 size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.5), Inches(6.60), Inches(12), Inches(0.35),
             "★ 首场目标：现场 80–120 人 · 直播 5,000+ 观看 · 3 篇媒体稿 · 200–400 人社群",
             size=11, bold=True, color=PLUM, italic=True,
             align=PP_ALIGN.CENTER)

    # ================= 12. 主视觉与传播 =================
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="07 · 传播",
               page_title="主视觉与传播设计",
               subtitle="东方极简美学 · 一滴水视觉母题 · 大留白")

    # 左：视觉规范
    add_text(s, Inches(0.5), Inches(1.15), Inches(6), Inches(0.4),
             "主视觉规范", size=14, bold=True, color=MOSS)
    add_rect(s, Inches(0.5), Inches(1.55), Inches(6), Inches(5.2),
             fill=IVORY, line=LINE)

    # 色板
    add_text(s, Inches(0.7), Inches(1.70), Inches(5.6), Inches(0.35),
             "色板（三色 + 素笺）", size=12, bold=True, color=MOSS)
    palettes = [("墨绿", MOSS), ("深梅", PLUM), ("金", GOLD), ("米白", IVORY)]
    for i, (n, c) in enumerate(palettes):
        x = Inches(0.7) + Inches(1.35) * i
        add_rect(s, x, Inches(2.10), Inches(1.15), Inches(0.6), fill=c,
                 line=LINE if c == IVORY else None)
        add_text(s, x, Inches(2.75), Inches(1.15), Inches(0.3), n,
                 size=10, bold=True, color=INK, align=PP_ALIGN.CENTER)

    # 字体
    add_text(s, Inches(0.7), Inches(3.20), Inches(5.6), Inches(0.35),
             "字体", size=12, bold=True, color=MOSS)
    add_text(s, Inches(0.7), Inches(3.55), Inches(5.6), Inches(1.0),
             "中文：宋体 / 思源宋体（人文气质）\n"
             "英文：Georgia / Playfair Display（衬线）\n"
             "标语字号：60pt+ · 正文：14–18pt · 注释：10–12pt",
             size=11, color=INK)

    # 版式原则
    add_text(s, Inches(0.7), Inches(4.60), Inches(5.6), Inches(0.35),
             "版式原则", size=12, bold=True, color=MOSS)
    principles = [
        "大留白（40%+）· 素笺气质",
        "主标题居中 · 副标题简洁",
        "一滴水涟漪母题（同心圆）",
        "黑白/单色人文摄影 · 忌花哨",
    ]
    for i, p in enumerate(principles):
        y = Inches(4.98) + Inches(0.36) * i
        add_round(s, Inches(0.75), y + Inches(0.08), Inches(0.20), Inches(0.20),
                  "●", fill=MOSS, color=IVORY, size=8, bold=True)
        add_text(s, Inches(1.05), y, Inches(5.2), Inches(0.35), p,
                 size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # 右：传播矩阵
    add_text(s, Inches(7.0), Inches(1.15), Inches(6), Inches(0.4),
             "传播矩阵", size=14, bold=True, color=MOSS)
    channels = [
        ("微信公众号", "三家主办方联合发稿 · 内容首发", MOSS),
        ("视频号 / 抖音", "3–5 支 60 秒短视频 · 长期沉淀", PLUM),
        ("微信社群", "预报名 + 长期社群 + 后续通告", GOLD),
        ("主流媒体", "澎湃 / 一财 / 21 世纪 / 界面 深度报道", DEEP),
        ("学术期刊 / 内刊", "思想内容长期沉淀 · 出版路径", MOSS),
        ("直播", "视频号直播 · 目标 5,000+ 观看", PLUM),
    ]
    for i, (t, d, c) in enumerate(channels):
        y = Inches(1.55) + Inches(0.85) * i
        add_rect(s, Inches(7.0), y, Inches(0.18), Inches(0.75), fill=c)
        add_rect(s, Inches(7.20), y, Inches(5.65), Inches(0.75),
                 fill=IVORY, line=LINE)
        add_text(s, Inches(7.35), y + Inches(0.05), Inches(2.0), Inches(0.32),
                 t, size=12, bold=True, color=MOSS, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(7.35), y + Inches(0.35), Inches(5.4), Inches(0.4),
                 d, size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # ================= 13. 组织架构与 RACI =================
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="07 · 组织",
               page_title="组织架构与 RACI 分工",
               subtitle="三方联合体 · 明确责任 · 高效协同")

    # 左：核心岗位
    add_text(s, Inches(0.5), Inches(1.15), Inches(6), Inches(0.4),
             "核心岗位（首场）", size=14, bold=True, color=MOSS)
    roles = [
        ("项目总召集人", "战略统筹 · 嘉宾锁定 · 三方协调", "首乾书院", MOSS),
        ("学术策划", "内容策划 · 主题设计 · 学术把关", "复旦住房中心", PLUM),
        ("商务对接", "企业家/产业界邀请 · 政府关系", "杨浦科企联", GOLD),
        ("现场执行", "场地 / 布置 / 接待 / 晚宴执行", "首乾书院", DEEP),
        ("传播/媒体", "视觉设计 / 稿件 / 直播 / 视频号", "首乾书院 + 联合", MOSS),
        ("财务/合规", "预算 / 结算 / 合同 / 法务审核", "首乾书院", PLUM),
    ]
    for i, (t, d, dept, c) in enumerate(roles):
        y = Inches(1.55) + Inches(0.85) * i
        add_rect(s, Inches(0.5), y, Inches(0.18), Inches(0.75), fill=c)
        add_rect(s, Inches(0.7), y, Inches(5.7), Inches(0.75),
                 fill=IVORY, line=LINE)
        add_text(s, Inches(0.85), y + Inches(0.05), Inches(1.9), Inches(0.32),
                 t, size=12, bold=True, color=MOSS, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(0.85), y + Inches(0.35), Inches(3.6), Inches(0.4),
                 d, size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_round(s, Inches(4.55), y + Inches(0.20), Inches(1.75), Inches(0.35),
                  dept, fill=c, color=MOSS if c == GOLD else IVORY,
                  size=10, bold=True)

    # 右：RACI 表
    add_text(s, Inches(7.0), Inches(1.15), Inches(6), Inches(0.4),
             "RACI 关键事项", size=14, bold=True, color=MOSS)
    header = ["关键事项", "首乾", "复旦", "杨浦"]
    rows = [
        ["战略统筹与嘉宾锁定", "R+A", "C", "C"],
        ["内容策划与学术把关", "C", "R+A", "I"],
        ["企业家邀请与政府对接", "C", "I", "R+A"],
        ["现场执行与晚宴运营", "R+A", "I", "C"],
        ["传播与媒体", "R+A", "C", "C"],
        ["财务预算与合同", "R+A", "C", "C"],
    ]
    add_table(s, Inches(7.0), Inches(1.55), Inches(6), Inches(4.5),
              header, rows,
              col_widths=[Inches(2.6), Inches(1.15), Inches(1.15), Inches(1.10)],
              header_size=11, body_size=10, body_align=PP_ALIGN.CENTER)
    add_text(s, Inches(7.0), Inches(6.15), Inches(6.0), Inches(0.5),
             "R = Responsible · A = Accountable · C = Consulted · I = Informed",
             size=10, italic=True, color=GREY)

    # ================= 14. 预算与投入产出 =================
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="07 · 预算",
               page_title="首场预算与投入产出",
               subtitle="14–21 万元 · 打造品牌资产 · 为长期铺基础")

    # 左：预算表
    add_text(s, Inches(0.5), Inches(1.15), Inches(6), Inches(0.4),
             "首场预算（策划口径 · 万元）", size=14, bold=True, color=MOSS)
    budget_header = ["科目", "预算范围", "备注"]
    budget_rows = [
        ["场地租赁（一滴水）", "4–6", "半日 + 晚宴"],
        ["晚宴餐饮", "3–5", "40–60 人私享"],
        ["主视觉与物料", "1–2", "海报/KV/宴请手册"],
        ["摄影摄像", "1.5", "全程记录 + 剪辑"],
        ["直播/传播", "1.5", "视频号直播 + 短视频"],
        ["嘉宾接待与礼品", "2–3", "车费 / 伴手礼"],
        ["其他杂项", "1–2", "应急预留"],
        ["合计", "14–21", "★"],
    ]
    add_table(s, Inches(0.5), Inches(1.55), Inches(6), Inches(4.5),
              budget_header, budget_rows,
              col_widths=[Inches(2.4), Inches(1.4), Inches(2.2)],
              header_size=11, body_size=10)

    # 右：投入产出
    add_text(s, Inches(7.0), Inches(1.15), Inches(6), Inches(0.4),
             "首场后 6 个月投入产出", size=14, bold=True, color=MOSS)
    outputs = [
        ("品牌资产", "形成「北外滩·一滴水」品牌记忆点", MOSS),
        ("社群资产", "核心社群 200–400 人\n（学者+企业家+文化）", PLUM),
        ("内容资产", "1 次深度内容\n→ 3 篇公众号 + 5 支短视频", GOLD),
        ("媒体资产", "1 次媒体深度报道\n+ 5–10 家转载", DEEP),
        ("未来价值", "为 Y2 全年 4 场铺基础\n3 年内形成年度品牌", MOSS),
    ]
    for i, (t, d, c) in enumerate(outputs):
        y = Inches(1.55) + Inches(0.90) * i
        add_rect(s, Inches(7.0), y, Inches(0.18), Inches(0.80), fill=c)
        add_rect(s, Inches(7.20), y, Inches(5.7), Inches(0.80),
                 fill=IVORY, line=LINE)
        add_text(s, Inches(7.35), y + Inches(0.05), Inches(2.0), Inches(0.35),
                 t, size=13, bold=True, color=MOSS, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(9.4), y + Inches(0.05), Inches(3.5), Inches(0.70),
                 d, size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # 底部宣示
    add_rect(s, Inches(0.5), Inches(6.30), Inches(12.3), Inches(0.55),
             fill=MOSS)
    add_text(s, Inches(0.7), Inches(6.32), Inches(12.0), Inches(0.50),
             "★ 首场是「品牌种子」投资 —— 3 年后单场活动商业价值将放大 5–10 倍",
             size=12, bold=True, color=GOLD, anchor=MSO_ANCHOR.MIDDLE)

    # ================= 15. 风险与合规 =================
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="07 · 风险",
               page_title="风险与合规清单",
               subtitle="事前预防 · 事中管控 · 事后复盘")

    risks = [
        ("R1 · 内容合规", "低", "中",
         "主题分享内容与议题双方书面确认；敏感表述法务预审",
         "M-30 前完成"),
        ("R2 · 场地合规", "低", "高",
         "与一滴水签正式租赁合同（含消防/安全条款）",
         "签约前"),
        ("R3 · 餐饮合规", "低", "中",
         "晚宴执行方需具备餐饮资质 + 食品安全承诺",
         "M-30 前锁定"),
        ("R4 · 知识产权", "中", "中",
         "分享内容录音/录像需嘉宾书面授权 + 首乾书院统一保管",
         "M-14 前完成"),
        ("R5 · 社群管理", "中", "低",
         "预报名信息按 PIPL 处理 · 建立信息删除 SOP",
         "常态化"),
        ("R6 · 直播事故", "低", "中",
         "双机位备份 + 5G 独立信号 + 应急方案",
         "M-7 演练"),
        ("R7 · 嘉宾变更", "中", "高",
         "备用嘉宾 / 视频致辞方案 + 24h 内替换",
         "常态化"),
        ("R8 · 报名超额", "中", "低",
         "分级候补名单 + 直播替代 + 下场优先席",
         "M-14 起"),
    ]
    header = ["#", "风险类别", "概率", "影响", "对冲手段", "启动时点"]
    add_table(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.3),
              header, risks,
              col_widths=[Inches(1.6), Inches(1.4), Inches(0.8), Inches(0.8),
                          Inches(5.7), Inches(2.0)],
              header_size=11, body_size=10)

    add_text(s, Inches(0.5), Inches(6.60), Inches(12), Inches(0.35),
             "风险委员会：三方各 1 人 · 首场前每周评审 · 长期每月复盘",
             size=11, italic=True, color=GREY, align=PP_ALIGN.CENTER)

    # ================= 16. 12 个月里程碑 =================
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="07 · 里程碑",
               page_title="12 个月品牌里程碑",
               subtitle="从首场筹备 → 第二场落地 · 品牌立根")

    months = [f"M{i+1}" for i in range(12)]
    chart_x = Inches(2.8); chart_y = Inches(1.2); chart_w = Inches(10.0)
    col_w = chart_w / 12.0
    for i, m in enumerate(months):
        x = chart_x + col_w * i
        add_rect(s, x, chart_y, col_w, Inches(0.35), fill=MOSS)
        add_text(s, x, chart_y, col_w, Inches(0.35), m,
                 size=10, bold=True, color=IVORY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    tasks = [
        ("立项 + 三方联合体",       0, 1, MOSS),
        ("场地档期确认",            0, 2, GOLD),
        ("内容策划 + 主题设计",     1, 2, PLUM),
        ("嘉宾邀请函（首轮）",     1, 2, DEEP),
        ("主视觉 + 海报 + 报名页", 2, 2, GOLD),
        ("预报名启动（社群沉淀）", 3, 2, MOSS),
        ("正式邀请 + 席位分配",   4, 1, PLUM),
        ("首场活动（10.31）",     5, 1, GOLD),
        ("媒体复盘 + 内容沉淀",   5, 2, DEEP),
        ("第二场策划启动",         7, 2, MOSS),
        ("第二场预报名",           9, 2, PLUM),
        ("第二场活动（12 月）",   11, 1, GOLD),
    ]
    row_h = Inches(0.36)
    for i, (lbl, start, dur, c) in enumerate(tasks):
        y = chart_y + Inches(0.4) + (row_h + Inches(0.05)) * i
        add_text(s, Inches(0.4), y, Inches(2.35), row_h, lbl,
                 size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, chart_x, y, chart_w, row_h, fill=IVORY, line=LINE)
        bx = chart_x + col_w * start
        bw = col_w * dur
        add_rect(s, bx + Emu(20000), y + Emu(20000),
                 bw - Emu(40000), row_h - Emu(40000), fill=c)

    # ================= 17. 品牌延伸 =================
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="07 · 延伸",
               page_title="品牌延伸与商业价值",
               subtitle="5 大延伸路径 · 长期商业化可能性")

    extensions = [
        ("图书出版", "思想内容整理成书\n长期沉淀 IP 资产",
         "与出版社联合出品\n目标：3 年 1 本", MOSS),
        ("音频/视频课程", "与得到/樊登/混沌合作\n输出精华课程",
         "首场后 6 个月上线\n目标：3 万人订阅", PLUM),
        ("企业内训", "私享晚宴模式\n复用到企业内训场景",
         "定制化服务\n单价 20–50 万/场", GOLD),
        ("文化空间联营", "与其他文化空间联动\n（西岸/思南公馆等）",
         "共建「上海人文思想地图」",
         DEEP),
        ("城市巡回", "北京/杭州/深圳\n姊妹活动",
         "Y3 启动\n目标：4 城联动",
         MOSS),
    ]
    ew = Inches(2.45); ey = Inches(1.20); eh = Inches(5.4); egap = Inches(0.10)
    for i, (t, d, biz, c) in enumerate(extensions):
        x = Inches(0.5) + (ew + egap) * i
        add_rect(s, x, ey, ew, eh, fill=IVORY, line=LINE)
        add_rect(s, x, ey, ew, Inches(0.9), fill=c)
        add_text(s, x, ey + Inches(0.20), ew, Inches(0.5), t,
                 size=15, bold=True,
                 color=MOSS if c == GOLD else IVORY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 描述
        add_text(s, x + Inches(0.15), ey + Inches(1.0), ew - Inches(0.3),
                 Inches(1.8), d, size=11, color=INK, align=PP_ALIGN.CENTER)
        add_rect(s, x + Inches(0.20), ey + Inches(3.0), ew - Inches(0.4),
                 Emu(15000), fill=GOLD)
        add_text(s, x + Inches(0.15), ey + Inches(3.2), ew - Inches(0.3),
                 Inches(1.5), biz, size=10, color=PLUM, italic=True,
                 align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.35),
             "★ 五大延伸路径共同支撑 3 年品牌收益：从纯投入 → 品牌资产 → 商业化收益",
             size=11, italic=True, color=GREY, align=PP_ALIGN.CENTER)

    # ================= 18. 决策请求 =================
    s = new_slide(prs)
    add_chrome(s, prs, page_no=0, phase_label="决策请求",
               page_title="给主办联合体的决策请求",
               subtitle="6 项关键授权 · 全部落定即可启动 10.31 筹备")

    decisions = [
        ("① 立项与三方协议",
         "签署三方联合主办协议 · 明确权益与分工",
         "T+7d 内"),
        ("② 首场预算授权",
         "首场 14–21 万元预算获三方书面确认",
         "T+14d 内"),
        ("③ 场地档期锁定",
         "10.31（周六）一滴水场地档期正式预订",
         "T+7d 内"),
        ("④ 主嘉宾策划授权",
         "授权首乾书院牵头首场主题嘉宾邀请",
         "T+7d 内"),
        ("⑤ 传播方案与视觉基调",
         "主视觉基调（墨绿 + 米白 + 金）三方联合确认",
         "T+21d 内"),
        ("⑥ 三年品牌路线图",
         "认可「Y1 首场 → Y2 4 场 → Y3 品牌矩阵」路线图",
         "T+21d 内"),
    ]
    for i, (t, d, when) in enumerate(decisions):
        y = Inches(1.20) + Inches(0.90) * i
        add_rect(s, Inches(0.5), y, Inches(3.5), Inches(0.80), fill=MOSS)
        add_text(s, Inches(0.5), y, Inches(3.5), Inches(0.80), t,
                 size=15, bold=True, color=GOLD,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(4.10), y, Inches(6.5), Inches(0.80),
                 fill=IVORY, line=LINE)
        add_text(s, Inches(4.30), y, Inches(6.3), Inches(0.80), d,
                 size=12, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_round(s, Inches(10.75), y + Inches(0.20), Inches(2.1),
                  Inches(0.40), when, fill=GOLD, color=MOSS,
                  size=12, bold=True)

    # 底部
    add_rect(s, Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.35),
             fill=PLUM)
    add_text(s, Inches(0.7), Inches(6.66), Inches(12.0), Inches(0.33),
             "★ 6 项授权全部落定后，10.31 筹备正式启动 · T-90 天倒计时开始",
             size=11, bold=True, color=IVORY, italic=True,
             anchor=MSO_ANCHOR.MIDDLE)

    # ================= 19. 联合宣言 =================
    s = new_slide(prs)
    add_rect(s, 0, 0, sw, sh, fill=IVORY)
    add_rect(s, 0, 0, sw, Emu(80000), fill=MOSS)
    add_rect(s, 0, sh - Emu(80000), sw, Emu(80000), fill=GOLD)

    # 一滴水（左侧）
    ripple_cx = Inches(2.5); ripple_cy = Inches(3.75)
    for r in [2.2, 1.7, 1.2, 0.7, 0.3]:
        rad = Inches(r)
        add_oval(s, ripple_cx - rad, ripple_cy - rad, rad * 2, rad * 2,
                 line_color=MOSS if r != 0.3 else PLUM,
                 line_w=1 if r > 0.5 else 2.5,
                 fill=None if r > 0.35 else PLUM)

    # 宣言文字（右侧）
    add_text(s, Inches(5.5), Inches(1.5), Inches(7.5), Inches(0.5),
             "联合宣言", size=13, bold=True, color=GOLD, italic=True)
    add_text(s, Inches(5.5), Inches(2.0), Inches(7.5), Inches(0.9),
             "北外滩 · 一滴水", size=32, bold=True, color=MOSS)
    add_text(s, Inches(5.5), Inches(2.85), Inches(7.5), Inches(0.9),
             "人文思想之夜", size=32, bold=True, color=PLUM)
    add_rect(s, Inches(5.5), Inches(3.90), Inches(2.0), Emu(30000), fill=GOLD)

    add_text(s, Inches(5.5), Inches(4.15), Inches(7.5), Inches(1.5),
             "以一滴水的谦逊\n"
             "承汪洋大海的深邃\n"
             "让思想在北外滩生根",
             size=17, color=INK, italic=True)

    add_text(s, Inches(5.5), Inches(5.90), Inches(7.5), Inches(0.35),
             "首场 · 2026 年 10 月 31 日 · 上海北外滩 · 一滴水",
             size=12, color=MOSS, bold=True)
    add_text(s, Inches(5.5), Inches(6.35), Inches(7.5), Inches(0.35),
             "复旦大学住房政策研究中心 · 首乾书院 · 上海市杨浦区科技企业联合会",
             size=10, color=GREY)

    # ================= 20. Q&A / 致谢 =================
    s = new_slide(prs)
    add_rect(s, 0, 0, sw, sh, fill=MOSS)
    add_rect(s, 0, Inches(3.6), sw, Emu(40000), fill=GOLD)

    # 一滴水（中央）
    ripple_cx = Inches(6.66); ripple_cy = Inches(3.75)
    for r in [3.2, 2.5, 1.8, 1.1, 0.5, 0.15]:
        rad = Inches(r)
        add_oval(s, ripple_cx - rad, ripple_cy - rad, rad * 2, rad * 2,
                 line_color=GOLD if r > 0.4 else PLUM,
                 line_w=0.8 if r > 0.7 else 2.5,
                 fill=None if r > 0.2 else PLUM)

    add_text(s, Inches(0.8), Inches(2.0), Inches(11), Inches(0.5),
             "BEIWAITAN · Humanities Nights", size=14, bold=True,
             color=GOLD, italic=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), Inches(4.5), Inches(11), Inches(0.7),
             "北外滩 · 一滴水", size=36, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), Inches(5.2), Inches(11), Inches(0.7),
             "人文思想之夜", size=36, bold=True, color=IVORY,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), Inches(6.15), Inches(11), Inches(0.5),
             "Q & A · 谢 谢", size=18, bold=True, color=IVORY,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), Inches(6.75), Inches(11), Inches(0.35),
             "主办联合体 · 品牌规划",
             size=11, color=IVORY, italic=True, align=PP_ALIGN.CENTER)

    # 回填页码
    total = len(SLIDES)
    for sl in SLIDES:
        for shape in sl.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text and run.text.strip().endswith(" / 0"):
                        run.text = run.text.replace(" / 0", f" / {total}")

    out = Path(__file__).resolve().parent.parent / "docs" / "deck" / \
        "北外滩人文思想之夜-品牌规划.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"✓ Deck written: {out}  ({total} slides)")
    return out


if __name__ == "__main__":
    build()
