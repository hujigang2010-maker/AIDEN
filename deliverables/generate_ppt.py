# -*- coding: utf-8 -*-
"""
生成《与夏春先生 · 长三角战略合作建议方案》PPT
机构：复旦大学住房政策研究中心 × 上海市杨浦区科技企业联合会 × 上海市科技企业联合会
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- 主题配色 ----------
NAVY = RGBColor(0x14, 0x2A, 0x4A)      # 深藏青 主色
NAVY2 = RGBColor(0x1F, 0x3A, 0x5F)     # 次藏青
GOLD = RGBColor(0xC9, 0xA2, 0x27)      # 金 强调
GOLD_L = RGBColor(0xE7, 0xD3, 0x8E)    # 浅金
GREY = RGBColor(0x5B, 0x64, 0x70)      # 正文灰
LGREY = RGBColor(0xEE, 0xF1, 0xF5)     # 浅底
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEAL = RGBColor(0x2E, 0x7D, 0x83)      # 辅助
CARD = RGBColor(0xF6, 0xF8, 0xFB)

FONT = "Microsoft YaHei"
FONT_L = "Microsoft YaHei Light"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

PAGE = {"n": 0}


def _set_font(run, size, color=NAVY, bold=False, name=FONT, italic=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = name
    # 中文字体
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', name)


def add_rect(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE, shadow=False):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    if shadow:
        el = sp._element.spPr
        effst = el.makeelement(qn('a:effectLst'), {})
        outer = el.makeelement(qn('a:outerShdw'),
                               {'blurRad': '60000', 'dist': '25000', 'dir': '5400000', 'rotWithShape': '0'})
        clr = el.makeelement(qn('a:srgbClr'), {'val': '9AA6B2'})
        alpha = el.makeelement(qn('a:alpha'), {'val': '55000'})
        clr.append(alpha)
        outer.append(clr)
        effst.append(outer)
        el.append(effst)
    return sp


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             wrap=True, space_after=6, line_spacing=1.05):
    """runs: list of paragraphs; each paragraph is list of (text,size,color,bold,name,italic)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for seg in para:
            text, size, color, bold = seg[0], seg[1], seg[2], seg[3]
            name = seg[4] if len(seg) > 4 else FONT
            italic = seg[5] if len(seg) > 5 else False
            r = p.add_run()
            r.text = text
            _set_font(r, size, color, bold, name, italic)
    return tb


def bg(slide, color=WHITE):
    add_rect(slide, 0, 0, SW, SH, color)


def footer(slide, title):
    PAGE["n"] += 1
    add_rect(slide, 0, SH - Inches(0.4), SW, Inches(0.4), NAVY)
    add_text(slide, Inches(0.5), SH - Inches(0.39), Inches(9), Inches(0.36),
             [[("复旦大学住房政策研究中心 × 上海市杨浦区科技企业联合会 × 上海市科技企业联合会", 9, GOLD_L, False)]],
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, SW - Inches(4.5), SH - Inches(0.39), Inches(4), Inches(0.36),
             [[(f"{title}    |    {PAGE['n']:02d}", 9, WHITE, False)]],
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def content_header(slide, kicker, title):
    add_rect(slide, 0, 0, SW, Inches(1.25), NAVY)
    add_rect(slide, 0, Inches(1.25), SW, Inches(0.06), GOLD)
    add_rect(slide, Inches(0.5), Inches(0.30), Inches(0.12), Inches(0.62), GOLD)
    add_text(slide, Inches(0.78), Inches(0.26), Inches(11.5), Inches(0.34),
             [[(kicker, 12, GOLD_L, False)]], anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(0.78), Inches(0.58), Inches(11.7), Inches(0.55),
             [[(title, 25, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)


def bullet_card(slide, x, y, w, h, icon_text, head, body, accent=GOLD):
    add_rect(slide, x, y, w, h, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    add_rect(slide, x, y, Inches(0.11), h, accent, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(slide, x + Inches(0.28), y + Inches(0.16), w - Inches(0.5), Inches(0.4),
             [[(icon_text + "  ", 14, accent, True), (head, 14.5, NAVY, True)]])
    add_text(slide, x + Inches(0.28), y + Inches(0.62), w - Inches(0.5), h - Inches(0.7),
             [[(body, 11.5, GREY, False)]], line_spacing=1.12)


# =========================================================
# Slide 1 — 封面
# =========================================================
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
add_rect(s, 0, 0, SW, SH, NAVY)
# 装饰
add_rect(s, 0, 0, Inches(0.28), SH, GOLD)
add_rect(s, SW - Inches(4.6), 0, Inches(4.6), SH, NAVY2)
add_rect(s, SW - Inches(4.6), 0, Inches(0.06), SH, GOLD)
add_text(s, Inches(0.9), Inches(1.35), Inches(8), Inches(0.5),
         [[("战略合作建议方案  |  STRATEGIC PARTNERSHIP PROPOSAL", 13, GOLD_L, False)]])
add_text(s, Inches(0.9), Inches(2.0), Inches(8.2), Inches(2.2),
         [[("与夏春先生", 46, WHITE, True)],
          [("长三角区域战略合作建议方案", 33, WHITE, True)]], line_spacing=1.05)
add_rect(s, Inches(0.95), Inches(4.15), Inches(2.6), Inches(0.05), GOLD)
add_text(s, Inches(0.9), Inches(4.4), Inches(8), Inches(1.0),
         [[("核心定位：", 15, GOLD_L, True), ("“夏春财经智识”长三角区域独家总代理", 15, WHITE, True)],
          [("交付载体：本方案（PPT）+ 配套测算与清单（Excel）", 12.5, RGBColor(0xC7,0xD2,0xE0), False)]],
         line_spacing=1.25)
add_text(s, Inches(0.9), Inches(6.35), Inches(8.5), Inches(0.8),
         [[("提报机构：复旦大学住房政策研究中心 × 上海市杨浦区科技企业联合会 × 上海市科技企业联合会", 12, WHITE, False)],
          [("2026 年 7 月", 11.5, GOLD_L, False)]], line_spacing=1.3)
# 右侧竖排要点
add_text(s, SW - Inches(4.3), Inches(1.4), Inches(3.6), Inches(4.8),
         [[("合作要点", 16, GOLD, True)],
          [("", 6, WHITE, False)],
          [("① 长期深化合作 + 3 个月适度过渡", 12.5, WHITE, False)],
          [("", 4, WHITE, False)],
          [("② 长三角区域服务总代理", 12.5, WHITE, False)],
          [("", 4, WHITE, False)],
          [("③ 依托复旦 + 市 / 区科技企业联合会背书", 12.5, WHITE, False)],
          [("", 4, WHITE, False)],
          [("④ 只做新区域增量，区隔新老团队", 12.5, WHITE, False)],
          [("", 4, WHITE, False)],
          [("⑤ 以夏春观点/语录为内容内核", 12.5, WHITE, False)]], line_spacing=1.15)

# =========================================================
# Slide 2 — 目录
# =========================================================
s = prs.slides.add_slide(BLANK)
bg(s)
content_header(s, "CONTENTS", "方案框架 · 目录")
items = [
    ("01", "合作背景与机遇"), ("02", "夏春先生 IP 与核心观点（语录）解析"),
    ("03", "我方机构与资源禀赋"), ("04", "合作定位：长三角区域总代理"),
    ("05", "差异化：区别于夏春新老团队"), ("06", "合作阶段：3 个月过渡 + 长期战略"),
    ("07", "长三角落地场景与内容矩阵"), ("08", "分工、协作与商业模式"),
    ("09", "关键问题与风险管控"), ("10", "三个月过渡期行动计划与下一步"),
]
cols, x0, y0 = 2, Inches(0.7), Inches(1.6)
cw, ch, gx, gy = Inches(5.9), Inches(0.82), Inches(0.15), Inches(0.2)
for i, (num, txt) in enumerate(items):
    r, c = divmod(i, cols)
    x = x0 + c * (cw + gx)
    y = y0 + r * (ch + gy)
    add_rect(s, x, y, cw, ch, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    add_rect(s, x, y, Inches(0.9), ch, NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, x, y, Inches(0.9), ch, [[(num, 22, GOLD, True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(1.1), y, cw - Inches(1.2), ch, [[(txt, 15, NAVY, True)]],
             anchor=MSO_ANCHOR.MIDDLE)
footer(s, "目录")

# =========================================================
# Slide 3 — 合作背景与机遇
# =========================================================
s = prs.slides.add_slide(BLANK)
bg(s)
content_header(s, "01  BACKGROUND & OPPORTUNITY", "合作背景与机遇")
cards = [
    ("◆", "夏春先生 · 稀缺财经 IP",
     "知名宏观经济学家，上善资本首席经济学家、中国首席经济学家论坛理事；北大本硕、明尼苏达经济学博士，曾任港大金融系教授。个人 IP“夏春财经智识”覆盖公众号 / 视频号 / B 站，内容含金量高、受众精准。", GOLD),
    ("◆", "内容矩阵高度契合高净值需求",
     "聚焦全球宏观、资产配置、投资策略与行为金融：美联储与美债美股、黄金、香港虚拟资产 / RWA、中国经济“韧性”与 K 型分化，正是长三角高净值与科创人群最关注的议题。", TEAL),
    ("◆", "长三角：增量最大的落地市场",
     "长三角高净值人群与产业资本高度密集，科创企业集群成熟。夏春线上影响力强、但区域线下承接与机构化运营长期缺位——这正是本次合作的价值缺口。", NAVY2),
    ("◆", "窗口已现 · 前期已充分沟通",
     "双方前期沟通已明确香港游学窗口期（中东资金、RWA 虚拟货币）、科技企业参访（AI / 科技前沿）等方向，且看重复旦 + 市 / 区科技企业联合会背书突破企业与高校资源——与我方禀赋精准匹配。", GOLD),
]
x0, y0, cw, ch, gx, gy = Inches(0.6), Inches(1.65), Inches(6.0), Inches(2.45), Inches(0.2), Inches(0.2)
for i, (ic, h, b, ac) in enumerate(cards):
    r, c = divmod(i, 2)
    bullet_card(s, x0 + c * (cw + gx), y0 + r * (ch + gy), cw, ch, ic, h, b, ac)
footer(s, "合作背景")

# =========================================================
# Slide 4 — 夏春语录解析（核心）
# =========================================================
s = prs.slides.add_slide(BLANK)
bg(s)
content_header(s, "02  KEY QUOTES & INSIGHTS", "夏春先生核心观点（语录）解析")
add_text(s, Inches(0.6), Inches(1.42), Inches(12), Inches(0.35),
         [[("说明：以下语录摘自夏春先生公开访谈 / 公众号 · 视频号内容，并结合双方前期沟通，提炼为长三角合作切入点。", 11, GREY, False, FONT, True)]])
quotes = [
    ("“2025 年……我个人更倾向于悲观派……需要更加警惕股债市场可能出现的危险。”",
     "稳健、风控优先的判断，天然契合长三角地产 / 高净值客群的避险配置焦虑。",
     "→ 打造“不确定市场下的资产防御”主题闭门沙龙，作为区域首发爆款。"),
    ("“每年年底做一张全年资产表现表：前一年涨幅大的减配、跌幅大的增配……做再平衡。”",
     "标志性、可复制的再平衡方法论，是可 IP 化的年度旗舰内容。",
     "→ 落地为长三角“年度资产配置盘点”旗舰活动 + 会员专栏，形成长期年度 IP。"),
    ("“香港对虚拟资产的接受度……速度甚至比美国更快；现货 ETF 允许实物兑换，全球首创。”",
     "直接呼应双方共同关注的香港游学 RWA / 中东资金窗口期。",
     "→ 我方作为长三角组团总入口，承接“半天讲座 + 半天参访”香港游学。"),
    ("“可将加密资产视为美股七巨头之外的‘第八个巨头’，与美股相关性最高约 0.5。”",
     "行为金融 + 另类资产的独到洞见，精准命中科创 / 家办人群。",
     "→ 面向杨浦科创企业家、家族办公室的另类配置主题专场。"),
    ("“我们更注重多元化分散、风险控制，寻找相对估值偏低但属同一投资逻辑的资产。”",
     "可迁移至不动产再配置，与复旦住房政策研究中心研究方向互补。",
     "→ 联合复旦推出“住房 / 不动产 + 大类资产”跨界研讨与课题。"),
    ("（前期沟通）“香港游学需夏春先生全程参与；方案由夏春先生主导决策。”",
     "夏春个人时间是核心稀缺资源，须为其减负。",
     "→ 我方承接全部执行、排期与履约，并以录播 / 图文复用降低其时间投入。"),
]
x0, y0, cw, ch, gx, gy = Inches(0.55), Inches(1.85), Inches(4.03), Inches(2.5), Inches(0.13), Inches(0.15)
for i, (q, interp, action) in enumerate(quotes):
    r, c = divmod(i, 3)
    x = x0 + c * (cw + gx)
    y = y0 + r * (ch + gy)
    add_rect(s, x, y, cw, ch, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    add_rect(s, x, y, cw, Inches(0.1), GOLD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, x + Inches(0.2), y + Inches(0.18), cw - Inches(0.4), Inches(1.0),
             [[(q, 10.3, NAVY, True)]], line_spacing=1.08)
    add_text(s, x + Inches(0.2), y + Inches(1.18), cw - Inches(0.4), Inches(0.7),
             [[("解读：", 9.5, TEAL, True), (interp, 9.5, GREY, False)]], line_spacing=1.05)
    add_text(s, x + Inches(0.2), y + Inches(1.86), cw - Inches(0.4), Inches(0.58),
             [[(action, 9.7, NAVY2, True)]], line_spacing=1.05)
footer(s, "语录解析")

# =========================================================
# Slide 5 — 我方机构与资源
# =========================================================
s = prs.slides.add_slide(BLANK)
bg(s)
content_header(s, "03  OUR INSTITUTIONS", "我方机构与资源禀赋")
# 左右两大机构卡片
add_rect(s, Inches(0.6), Inches(1.7), Inches(5.9), Inches(3.4), CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
add_rect(s, Inches(0.6), Inches(1.7), Inches(5.9), Inches(0.72), NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(0.85), Inches(1.7), Inches(5.5), Inches(0.72),
         [[("复旦大学住房政策研究中心", 17, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(0.9), Inches(2.6), Inches(5.3), Inches(2.4),
         [[("• 顶尖高校学术背书与公信力，提升活动 / 内容权威性", 12, GREY, False)],
          [("• 住房与不动产政策研究专长，与夏春大类资产配置互补", 12, GREY, False)],
          [("• 政、企、学高端资源网络与课题合作能力", 12, GREY, False)],
          [("• 承接高规格闭门研讨、政企对话与学术型内容", 12, GREY, False)]], line_spacing=1.35)
add_rect(s, Inches(6.8), Inches(1.7), Inches(5.9), Inches(3.4), CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
add_rect(s, Inches(6.8), Inches(1.7), Inches(5.9), Inches(0.72), TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(7.05), Inches(1.7), Inches(5.5), Inches(0.72),
         [[("上海市杨浦区科技企业联合会（为主）", 15, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(7.1), Inches(2.6), Inches(5.3), Inches(2.4),
         [[("• 杨浦（大创智 / 在线新经济）科创企业集群与会员企业", 12, GREY, False)],
          [("• 科技企业参访的本地落地与组织能力", 12, GREY, False)],
          [("• 联合上海市科技企业联合会，市 / 区两级协同、资源更广", 12, GREY, False)],
          [("• 承接 AI / 科技前沿主题参访、企业内训与产业沙龙", 12, GREY, False)]], line_spacing=1.35)
# 组合结论条
add_rect(s, Inches(0.6), Inches(5.35), Inches(12.1), Inches(1.4), NAVY2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_rect(s, Inches(0.6), Inches(5.35), Inches(0.13), Inches(1.4), GOLD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(0.95), Inches(5.5), Inches(11.6), Inches(1.15),
         [[("组合优势 = 承接夏春 IP 的最佳区域平台", 15, GOLD_L, True)],
          [("学术公信力（复旦） + 科创产业落地（市 / 区科技企业联合会） + 长三角本地网络，一站式解决夏春线下承接与机构化运营缺口，", 12, WHITE, False)],
          [("既能做高端内容与政企研讨，又能规模化落地企业参访、财富沙龙与出海游学。", 12, WHITE, False)]],
         line_spacing=1.2)
footer(s, "机构资源")

# =========================================================
# Slide 6 — 合作定位：总代理
# =========================================================
s = prs.slides.add_slide(BLANK)
bg(s)
content_header(s, "04  POSITIONING", "合作定位：长三角区域总代理")
add_rect(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(1.15), NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
add_text(s, Inches(0.95), Inches(1.6), Inches(11.4), Inches(1.15),
         [[("核心目标：成为“夏春财经智识”在长三角区域相关服务的独家 / 优先总代理", 17, WHITE, True)]],
         anchor=MSO_ANCHOR.MIDDLE)
subs = [
    ("服务范围", "区域线下活动、闭门会 / 研讨、科技企业参访、财富与资产配置沙龙、企业内训、内容本地化与私域运营，以及香港 / 出海游学的长三角组团总入口。", GOLD),
    ("授权模式", "区域独家授权 + 品牌联合（“夏春财经智识 · 长三角”）+ 收益分成；夏春侧保留 IP 与内容所有权，我方获区域运营与获客授权。", TEAL),
    ("代理价值", "夏春一个区域出口、一套结算、一个履约团队；把稀缺的个人时间从繁杂区域事务中解放，专注内容与关键出席。", NAVY2),
]
x0, y0, cw, ch, gy = Inches(0.6), Inches(2.95), Inches(12.1), Inches(1.15), Inches(0.18)
for i, (h, b, ac) in enumerate(subs):
    y = y0 + i * (ch + gy)
    add_rect(s, x0, y, cw, ch, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    add_rect(s, x0, y, Inches(2.3), ch, ac, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, x0, y, Inches(2.3), ch, [[(h, 15, WHITE, True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x0 + Inches(2.55), y, cw - Inches(2.8), ch,
             [[(b, 12, GREY, False)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
footer(s, "合作定位")

# =========================================================
# Slide 7 — 差异化：区别于新老团队
# =========================================================
s = prs.slides.add_slide(BLANK)
bg(s)
content_header(s, "05  DIFFERENTIATION", "差异化定位：区别于夏春新老团队")
rows = [
    ("对比维度", "夏春 · 老团队", "夏春 · 新团队", "我方（长三角总代理）"),
    ("核心侧重", "内容生产、全国泛在影响", "金融产品与新业务孵化", "长三角区域落地与机构化承接"),
    ("目标客群", "全国粉丝 / 线上受众", "金融端存量客户", "长三角高净值 + 科创企业 + 高校 / 政府"),
    ("主要场景", "视频号 / 公众号内容", "金融配置与产品", "区域活动 / 参访 / 游学组团 / 内训"),
    ("背书资源", "夏春个人 IP", "资本 / 金融资源", "复旦 + 市 / 区科技企业联合会 背书 + 本地网络"),
    ("客户归属", "既有存量", "金融客户", "只做长三角新增增量，签不重叠条款"),
]
tx, ty = Inches(0.6), Inches(1.75)
tw = Inches(12.13)
col_w = [Inches(2.03), Inches(3.0), Inches(3.0), Inches(4.1)]
row_h = Inches(0.78)
x = tx
for j in range(4):
    y = ty
    for i in range(len(rows)):
        cell_w = col_w[j]
        if i == 0:
            fill = NAVY
        elif j == 3:
            fill = RGBColor(0xEB, 0xE2, 0xC4)  # 高亮我方列
        else:
            fill = CARD if i % 2 else LGREY
        add_rect(s, x, y, cell_w, row_h, fill, line=WHITE)
        is_head = i == 0
        color = WHITE if is_head else (NAVY if (j == 0 or j == 3) else GREY)
        bold = is_head or j == 0 or j == 3
        add_text(s, x + Inches(0.12), y, cell_w - Inches(0.24), row_h,
                 [[(rows[i][j], 11.5 if is_head else 11, color, bold)]],
                 align=(PP_ALIGN.CENTER if j == 0 or is_head else PP_ALIGN.LEFT),
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        y += row_h
    x += col_w[j]
add_text(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.5),
         [[("原则：不争夺存量、聚焦区域增量；边界清晰，与新老团队互补而非重叠（明确主导权 / 客户归属）。", 11.5, NAVY2, True)]])
footer(s, "差异化")

# =========================================================
# Slide 8 — 合作阶段
# =========================================================
s = prs.slides.add_slide(BLANK)
bg(s)
content_header(s, "06  PHASING", "合作阶段：3 个月过渡 + 长期战略")
# 短期
add_rect(s, Inches(0.6), Inches(1.7), Inches(5.9), Inches(3.9), CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
add_rect(s, Inches(0.6), Inches(1.7), Inches(5.9), Inches(0.8), TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(0.85), Inches(1.7), Inches(5.5), Inches(0.8),
         [[("适度过渡 · 约 3 个月（双方适应与适配）", 15, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(0.9), Inches(2.7), Inches(5.4), Inches(2.8),
         [[("• 双方适应与适配，建立信任与协作机制", 12, GREY, False)],
          [("• 小步落地：单场跑通盈利模型，控风险", 12, GREY, False)],
          [("• 界定客户归属与收益分配规则", 12, GREY, False)],
          [("• 签署合作备忘录（MOU）与区域独家授权草案", 12, GREY, False)],
          [("• 敏感事项分步磨合，节奏可控", 12, GREY, False)]], line_spacing=1.5)
# 长期
add_rect(s, Inches(6.8), Inches(1.7), Inches(5.9), Inches(3.9), CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
add_rect(s, Inches(6.8), Inches(1.7), Inches(5.9), Inches(0.8), NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(7.05), Inches(1.7), Inches(5.5), Inches(0.8),
         [[("长期合作 · 深度战略协作", 15, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(7.1), Inches(2.7), Inches(5.4), Inches(2.8),
         [[("• 区域独家总代理授权（多年期）", 12, GREY, False)],
          [("• 收益分配清晰、数据共享，持续深化合作", 12, GREY, False)],
          [("• 共建长三角会员生态与 CRM 系统", 12, GREY, False)],
          [("• 年度活动 / 内容 / 游学 IP 化与品牌联合", 12, GREY, False)],
          [("• 由点及面复制至长三角多城市", 12, GREY, False)]], line_spacing=1.5)
# gate
add_rect(s, Inches(0.6), Inches(5.85), Inches(12.1), Inches(0.9), GOLD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(0.9), Inches(5.85), Inches(11.5), Inches(0.9),
         [[("过渡与转长期安排：", 13, NAVY, True),
           ("三个月为双方适应与适配的适度过渡期；期间运作顺畅、无重大问题，即自动转为长期区域独家总代理，无需另设考核门槛。", 12.5, NAVY, False)]],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
footer(s, "合作阶段")

# =========================================================
# Slide 9 — 落地场景与内容矩阵
# =========================================================
s = prs.slides.add_slide(BLANK)
bg(s)
content_header(s, "07  SCENARIOS & CONTENT MATRIX", "长三角落地场景与内容矩阵")
cards = [
    ("🏢", "科技企业参访", "AI / 科技前沿主题，依托杨浦科创集群 + 复旦，月度 1–2 场；结合夏春“七巨头 / 第八巨头”科技视角。", GOLD),
    ("📈", "财富与资产配置沙龙", "夏春年度资产表 / 再平衡框架、黄金、美债美股；面向高净值客群的闭门研讨。", TEAL),
    ("✈", "香港 / 出海游学组团", "RWA、虚拟资产、中东资金窗口；半天讲座 + 半天参访，我方作长三角组团总入口。", NAVY2),
    ("📱", "内容本地化与私域", "公众号 / 视频号“长三角专栏”、会员分层与社群运营，沉淀区域私域资产。", GOLD),
    ("🎓", "企业内训与政企研讨", "复旦住房政策中心 + 市 / 区科技企业联合会，承接企业内训与高规格政企对话。", TEAL),
    ("🧭", "家办 / 另类配置专场", "面向科创企业家、家族办公室的另类资产与行为金融专题。", NAVY2),
]
x0, y0, cw, ch, gx, gy = Inches(0.6), Inches(1.65), Inches(4.0), Inches(2.4), Inches(0.15), Inches(0.15)
for i, (ic, h, b, ac) in enumerate(cards):
    r, c = divmod(i, 3)
    bullet_card(s, x0 + c * (cw + gx), y0 + r * (ch + gy), cw, ch, ic, h, b, ac)
footer(s, "落地矩阵")

# =========================================================
# Slide 10 — 分工、协作与商业模式
# =========================================================
s = prs.slides.add_slide(BLANK)
bg(s)
content_header(s, "08  ROLES & BUSINESS MODEL", "分工、协作与商业模式")
# 分工两栏
add_rect(s, Inches(0.6), Inches(1.65), Inches(6.0), Inches(2.55), CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
add_text(s, Inches(0.85), Inches(1.8), Inches(5.6), Inches(0.5), [[("夏春侧", 15, GOLD, True)]])
add_text(s, Inches(0.9), Inches(2.35), Inches(5.5), Inches(1.8),
         [[("• IP、核心观点与讲座 / 关键出席", 11.5, GREY, False)],
          [("• 品牌授权与背书", 11.5, GREY, False)],
          [("• 关键嘉宾 / 高端资源引荐", 11.5, GREY, False)]], line_spacing=1.35)
add_rect(s, Inches(6.75), Inches(1.65), Inches(6.0), Inches(2.55), CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
add_text(s, Inches(7.0), Inches(1.8), Inches(5.6), Inches(0.5), [[("我方（总代理）", 15, TEAL, True)]])
add_text(s, Inches(7.05), Inches(2.35), Inches(5.5), Inches(1.8),
         [[("• 区域获客、场地与全流程执行", 11.5, GREY, False)],
          [("• 企业 / 高校 / 政府对接与背书", 11.5, GREY, False)],
          [("• 私域、CRM、组团与履约结算", 11.5, GREY, False)]], line_spacing=1.35)
# 商业模式条
add_rect(s, Inches(0.6), Inches(4.4), Inches(12.15), Inches(2.2), NAVY2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
add_rect(s, Inches(0.6), Inches(4.4), Inches(0.13), Inches(2.2), GOLD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(0.95), Inches(4.55), Inches(11.6), Inches(2.0),
         [[("商业模式与结算", 15, GOLD_L, True)],
          [("收入来源：活动票务 / 会员费 / 企业赞助与内训 / 游学组团佣金 / 内容订阅。", 12, WHITE, False)],
          [("分润逻辑：按角色与投入分层结算；成本相对固定，毛利率高，核心取决于客单价与销量（详见附 Excel 测算）。", 12, WHITE, False)],
          [("结算方式：长三角区域收入统一归集 → 按约定向夏春侧结算授权费 / 分成，账目透明、按月对账。", 12, WHITE, False)]],
         line_spacing=1.3)
footer(s, "分工与商业模式")

# =========================================================
# Slide 11 — 关键问题与风险管控
# =========================================================
s = prs.slides.add_slide(BLANK)
bg(s)
content_header(s, "09  RISK & GOVERNANCE", "关键问题与风险管控")
risks = [
    ("主导权与客户归属", "明确长三角新增客户归属规则与交叉销售（房产 / 保险 / 基金）权益划分，写入 MOU。", GOLD),
    ("IP 档期与精力稀缺", "我方承接全部执行、批量化排期；以录播 / 图文复用降低夏春时间投入，关键场次才需本人出席。", TEAL),
    ("与新老团队边界", "只做长三角区域增量，签署互不重叠 / 不竞争条款，定期同步避免撞单。", NAVY2),
    ("合规与免责", "金融观点仅供参考、不构成投资建议；活动内容与产品销售分离，沿用夏春一贯免责声明。", GOLD),
]
x0, y0, cw, ch, gx, gy = Inches(0.6), Inches(1.7), Inches(6.0), Inches(2.45), Inches(0.2), Inches(0.2)
for i, (h, b, ac) in enumerate(risks):
    r, c = divmod(i, 2)
    bullet_card(s, x0 + c * (cw + gx), y0 + r * (ch + gy), cw, ch, "▲", h, b, ac)
add_text(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.5),
         [[("以上关键问题为双方前期沟通中的待解决事项，将在三个月适度过渡期内逐项落定并写入合作备忘录。", 11.5, NAVY2, True)]])
footer(s, "风险管控")

# =========================================================
# Slide 12 — 三个月行动计划与下一步
# =========================================================
s = prs.slides.add_slide(BLANK)
bg(s)
content_header(s, "10  90-DAY PLAN & NEXT STEPS", "三个月适度过渡期行动计划与下一步")
months = [
    ("第 1 月", "启动与共识", ["签署 MOU + 区域独家授权草案", "盘客 / 会员画像梳理", "首场活动策划立项"], TEAL),
    ("第 2 月", "首战落地", ["首场科技参访 / 资产沙龙落地", "香港游学预售启动", "私域 & CRM 搭建"], GOLD),
    ("第 3 月", "复盘与转长期", ["第 2–3 场活动落地", "运作复盘与双方磨合评估", "过渡顺畅 → 自动转长期"], NAVY),
]
x0, y0, cw, ch, gx = Inches(0.6), Inches(1.7), Inches(3.9), Inches(3.5), Inches(0.22)
for i, (m, theme, tasks, ac) in enumerate(months):
    x = x0 + i * (cw + gx)
    add_rect(s, x, y0, cw, ch, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    add_rect(s, x, y0, cw, Inches(0.95), ac, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, x, y0 + Inches(0.1), cw, Inches(0.5), [[(m, 18, WHITE, True)]],
             align=PP_ALIGN.CENTER)
    add_text(s, x, y0 + Inches(0.52), cw, Inches(0.4), [[(theme, 13, GOLD_L if ac != GOLD else WHITE, True)]],
             align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.3), y0 + Inches(1.2), cw - Inches(0.55), Inches(2.2),
             [[("• " + t, 12, GREY, False)] for t in tasks], line_spacing=1.4, space_after=8)
# 下一步条
add_rect(s, Inches(0.6), Inches(5.5), Inches(12.15), Inches(1.2), NAVY2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_rect(s, Inches(0.6), Inches(5.5), Inches(0.13), Inches(1.2), GOLD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(0.95), Inches(5.62), Inches(11.6), Inches(1.0),
         [[("下一步", 14, GOLD_L, True),
           ("：确认过渡期独家授权范围与首场活动档期 → 组建联合微信群同步进度 → 会签合作备忘录。", 13, WHITE, False)],
          [("交付：本 PPT 方案 + 配套 Excel（语录解析表、3 个月计划、活动矩阵、财务测算、风险待办、KPI）。", 12, RGBColor(0xC7,0xD2,0xE0), False)]],
         line_spacing=1.3)
footer(s, "行动计划")

# =========================================================
# Slide 13 — 封底
# =========================================================
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, 0, Inches(2.6), SW, Inches(0.06), GOLD)
add_text(s, Inches(0), Inches(2.9), SW, Inches(1.2),
         [[("期待与夏春先生携手，", 30, WHITE, True)],
          [("共建长三角财经内容与活动新生态", 30, GOLD_L, True)]],
         align=PP_ALIGN.CENTER, line_spacing=1.15)
add_text(s, Inches(0), Inches(5.2), SW, Inches(0.8),
         [[("复旦大学住房政策研究中心 × 上海市杨浦区科技企业联合会 × 上海市科技企业联合会", 13, WHITE, False)],
          [("2026 年 7 月", 12, GOLD_L, False)]],
         align=PP_ALIGN.CENTER, line_spacing=1.4)

out = "/workspace/deliverables/夏春长三角战略合作建议方案.pptx"
prs.save(out)
print("saved:", out, "slides:", len(prs.slides._sldIdLst))
