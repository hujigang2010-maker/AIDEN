# -*- coding: utf-8 -*-
"""生成《赴宁波港及宁波经济技术开发区考察交流策划案》PPT。"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

import content as C

# 主题：港湾深蓝 + 青石绿（避开紫 / 奶油 / 报纸风）
NAVY = RGBColor(0x0B, 0x3D, 0x5C)
DEEP = RGBColor(0x07, 0x2A, 0x40)
TEAL = RGBColor(0x1A, 0x6B, 0x6B)
SEA = RGBColor(0x2A, 0x8F, 0x8F)
AMBER = RGBColor(0xC4, 0x7E, 0x1A)
LIGHT = RGBColor(0xE8, 0xF0, 0xF4)
MIST = RGBColor(0xF3, 0xF7, 0xF9)
GREY = RGBColor(0x5A, 0x63, 0x6E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x22, 0x2A, 0x33)

FONT = "微软雅黑"
SW = Inches(13.333)
SH = Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def set_cjk(run, font=FONT):
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", font)


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, fill=None, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    return sp


def textbox(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    for i, item in enumerate(lines):
        text, size, color, bold = item[0], item[1], item[2], item[3]
        sa = item[4] if len(item) > 4 else 4
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(sa)
        p.line_spacing = 1.1
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        set_cjk(r)
    return tb


def header(slide, tag, title, subtitle=None):
    rect(slide, 0, 0, SW, Inches(1.12), fill=NAVY)
    rect(slide, 0, Inches(1.12), SW, Pt(3), fill=TEAL)
    textbox(slide, Inches(0.5), Inches(0.12), Inches(8), Inches(0.35),
            [(tag, 12, AMBER, True)])
    ln = [(title, 24, WHITE, True)]
    if subtitle:
        ln.append((subtitle, 12, RGBColor(0xB8, 0xCF, 0xDC), False))
    textbox(slide, Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.7), ln)


def footer(slide, page):
    textbox(slide, Inches(0.5), Inches(7.08), Inches(10), Inches(0.3),
            [("复旦大学住房政策研究中心 · 上海市杨浦区科技企业联合会 → 宁波经开区投促局",
              9, GREY, False)])
    textbox(slide, Inches(11.9), Inches(7.08), Inches(1.1), Inches(0.3),
            [(str(page), 9, GREY, False)], align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, title, lines, accent=TEAL):
    rect(slide, x, y, w, h, fill=MIST, line=RGBColor(0xD0, 0xDE, 0xE6))
    rect(slide, x, y, Pt(5), h, fill=accent)
    textbox(slide, x + Inches(0.2), y + Inches(0.12), w - Inches(0.3), Inches(0.35),
            [(title, 13, NAVY, True)])
    body = [(ln, 11, DARK, False, 5) for ln in lines]
    textbox(slide, x + Inches(0.2), y + Inches(0.48), w - Inches(0.3), h - Inches(0.55),
            body)


# ================================================================ 1. 封面
s = add_slide()
rect(s, 0, 0, SW, SH, fill=DEEP)
# 氛围条带（非卡片堆叠）
rect(s, 0, 0, Inches(0.18), SH, fill=TEAL)
rect(s, 0, Inches(4.85), SW, Inches(2.65), fill=NAVY)
textbox(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.4),
        [("考察交流策划案  ·  " + C.VERSION, 14, SEA, True)])
textbox(s, Inches(0.8), Inches(2.1), Inches(11.5), Inches(1.2),
        [(C.PROJECT_TITLE, 32, WHITE, True, 8)])
textbox(s, Inches(0.8), Inches(3.4), Inches(11.5), Inches(0.5),
        [(C.PROJECT_SUBTITLE, 16, AMBER, False)])
textbox(s, Inches(0.8), Inches(5.15), Inches(11.5), Inches(1.8),
        [(C.ORGANIZER_LINE, 14, WHITE, True, 8),
         ("提交对象：" + C.RECIPIENT, 14, SEA, True, 8),
         (C.PARTICIPANT_LINE, 13, RGBColor(0xB8, 0xCF, 0xDC), False, 8),
         (C.TIME_WINDOW + "  ·  " + C.DOC_DATE, 12, RGBColor(0x9A, 0xB8, 0xC8), False)],
        align=PP_ALIGN.LEFT)

# ================================================================ 2. 目录
s = add_slide()
header(s, "目录", "策划案结构", "提交贵局审阅 · 便于对接落地")
items = [
    ("01", "背景与目的"),
    ("02", "组织与人员"),
    ("03", "合作主题方向"),
    ("04", "宁波港资源筛选"),
    ("05", "互动链接与行程"),
    ("06", "成果与恳请支持"),
]
for i, (no, name) in enumerate(items):
    col = i % 3
    row = i // 3
    x = Inches(0.6) + col * Inches(4.1)
    y = Inches(2.0) + row * Inches(2.0)
    rect(s, x, y, Inches(3.8), Inches(1.6), fill=MIST)
    rect(s, x, y, Inches(3.8), Pt(4), fill=TEAL)
    textbox(s, x + Inches(0.25), y + Inches(0.35), Inches(3.3), Inches(0.5),
            [(no, 22, TEAL, True)])
    textbox(s, x + Inches(0.25), y + Inches(0.9), Inches(3.3), Inches(0.45),
            [(name, 16, NAVY, True)])
footer(s, 2)

# ================================================================ 3. 背景与目的
s = add_slide()
header(s, "01 背景与目的", "为何赴宁波 · 希望达成什么", "以一次考察开启沪甬长效投促合作")
textbox(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.2),
        [("项目背景", 14, TEAL, True, 8)] +
        [(f"· {b}", 12, DARK, False, 10) for b in C.BACKGROUND])
rect(s, Inches(6.8), Inches(1.4), Inches(5.9), Inches(5.2), fill=NAVY)
textbox(s, Inches(7.1), Inches(1.65), Inches(5.4), Inches(0.4),
        [("考察目的", 14, AMBER, True)])
textbox(s, Inches(7.1), Inches(2.2), Inches(5.4), Inches(4.1),
        [(f"{i}. {p}", 13, WHITE, False, 12) for i, p in enumerate(C.PURPOSE, 1)])
footer(s, 3)

# ================================================================ 4. 三方角色
s = add_slide()
header(s, "02 组织与接待", "三方分工清晰 · 便于落地执行", "")
w = Inches(3.9)
gap = Inches(0.2)
x0 = Inches(0.5)
for i, party in enumerate(C.PARTIES):
    x = x0 + (w + gap) * i
    accent = [TEAL, SEA, AMBER][i]
    card(s, x, Inches(1.5), w, Inches(5.1),
         party["role"],
         [party["name"], ""] + [f"· {d}" for d in party["duties"]],
         accent=accent)
footer(s, 4)

# ================================================================ 5. 人员
s = add_slide()
header(s, "03 考察人员", "复旦校友 + 上海科技企业家", C.SCALE)
for i, (cat, desc) in enumerate(C.PARTICIPANT_PROFILE):
    y = Inches(1.45) + i * Inches(1.15)
    rect(s, Inches(0.5), y, Inches(12.3), Inches(1.0), fill=MIST)
    rect(s, Inches(0.5), y, Inches(0.12), Inches(1.0), fill=TEAL if i % 2 == 0 else AMBER)
    textbox(s, Inches(0.9), y + Inches(0.18), Inches(3.2), Inches(0.65),
            [(cat, 15, NAVY, True)], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(4.3), y + Inches(0.18), Inches(8.2), Inches(0.65),
            [(desc, 13, DARK, False)], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 5)

# ================================================================ 6. 合作方向
s = add_slide()
header(s, "04 合作方向", "参访交流 + 长期投洽事项", "招商引资 · 科技孵化 · 产业导入 · PE/VC · CVC")
positions = [
    (Inches(0.45), Inches(1.4)),
    (Inches(6.85), Inches(1.4)),
    (Inches(0.45), Inches(4.15)),
    (Inches(6.85), Inches(4.15)),
]
accents = [TEAL, SEA, AMBER, NAVY]
for i, theme in enumerate(C.COOP_THEMES):
    x, y = positions[i]
    card(s, x, y, Inches(5.95), Inches(2.5),
         theme["name"],
         [f"· {p}" for p in theme["points"]],
         accent=accents[i])
footer(s, 6)

# ================================================================ 7. 筛选原则
s = add_slide()
header(s, "04 资源筛选", "宁波港不宜走马观花 · 按三原则精选", C.VALUE_CHAIN)
textbox(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.7),
        [(C.SCREEN_INTRO, 12, DARK, False)])
for i, (name, desc) in enumerate(C.SCREEN_CRITERIA):
    x = Inches(0.5) + (i % 4) * Inches(3.15)
    y = Inches(2.4) + (i // 4) * Inches(2.4)
    rect(s, x, y, Inches(3.0), Inches(3.5), fill=MIST)
    rect(s, x, y, Inches(3.0), Pt(5), fill=TEAL if i < 2 else AMBER)
    textbox(s, x + Inches(0.15), y + Inches(0.3), Inches(2.7), Inches(0.5),
            [(f"0{i+1}", 18, TEAL, True)])
    textbox(s, x + Inches(0.15), y + Inches(0.9), Inches(2.7), Inches(0.5),
            [(name, 16, NAVY, True)])
    textbox(s, x + Inches(0.15), y + Inches(1.55), Inches(2.7), Inches(1.6),
            [(desc, 13, DARK, False)])
footer(s, 7)

# ================================================================ 8. 港口线
s = add_slide()
header(s, "04 港口线精选", "核心必看 + 推荐加深 + 文化选配", "穿山港区 · 航运展示厅 · 梅山港区 · 港口博物馆")
for i, site in enumerate(C.PORT_SITES):
    col = i % 2
    row = i // 2
    x = Inches(0.4) + col * Inches(6.4)
    y = Inches(1.35) + row * Inches(2.7)
    rect(s, x, y, Inches(6.2), Inches(2.5), fill=MIST)
    rect(s, x, y, Pt(5), Inches(2.5), fill=TEAL if "核心" in site["tier"] else AMBER)
    textbox(s, x + Inches(0.2), y + Inches(0.12), Inches(5.8), Inches(0.35),
            [(f"【{site['tier']}】", 11, TEAL, True)])
    textbox(s, x + Inches(0.2), y + Inches(0.42), Inches(5.8), Inches(0.55),
            [(site["name"], 13, NAVY, True)])
    textbox(s, x + Inches(0.2), y + Inches(1.05), Inches(5.8), Inches(1.3),
            [(f"看：{site['see']}", 11, DARK, False, 4),
             (f"互动：{site['interact']}", 11, GREY, False)])
footer(s, 8)

# ================================================================ 9. 产业孵化线
s = add_slide()
header(s, "04 产业与孵化线", "智造标杆 + 新能源 + 硬科技 + 落地载体", "与上海企业家、投资机构需求精准匹配")
for i, site in enumerate(C.INDUSTRY_SITES):
    col = i % 2
    row = i // 2
    x = Inches(0.4) + col * Inches(6.4)
    y = Inches(1.35) + row * Inches(2.7)
    rect(s, x, y, Inches(6.2), Inches(2.5), fill=MIST)
    rect(s, x, y, Pt(5), Inches(2.5), fill=SEA if "核心" in site["tier"] else AMBER)
    textbox(s, x + Inches(0.2), y + Inches(0.12), Inches(5.8), Inches(0.3),
            [(f"【{site['tier']}】", 11, SEA, True)])
    textbox(s, x + Inches(0.2), y + Inches(0.4), Inches(5.8), Inches(0.5),
            [(site["name"], 13, NAVY, True)])
    textbox(s, x + Inches(0.2), y + Inches(1.0), Inches(5.8), Inches(1.3),
            [(f"链接：{site['link']}", 11, DARK, False, 4),
             (f"互动：{site['interact']}", 11, GREY, False)])
footer(s, 9)

# ================================================================ 10. 互动链接
s = add_slide()
header(s, "05 互动链接", "把参访变成沟通与转化", "看—问—对—跟 闭环")
for i, link in enumerate(C.INTERACTION_LINKS):
    col = i % 2
    row = i // 2
    x = Inches(0.4) + col * Inches(6.4)
    y = Inches(1.35) + row * Inches(2.7)
    rect(s, x, y, Inches(6.2), Inches(2.5), fill=DEEP if i % 2 else NAVY)
    textbox(s, x + Inches(0.25), y + Inches(0.2), Inches(5.7), Inches(0.4),
            [(link["name"], 14, AMBER, True)])
    textbox(s, x + Inches(0.25), y + Inches(0.65), Inches(5.7), Inches(0.35),
            [(f"{link['anchor']} ｜ {link['form']}", 11, SEA, False)])
    textbox(s, x + Inches(0.25), y + Inches(1.1), Inches(5.7), Inches(1.2),
            [(f"· {a}", 12, WHITE, False, 5) for a in link["agenda"]])
footer(s, 10)

# ================================================================ 11. 时间
s = add_slide()
header(s, "05 时间安排", C.TIME_WINDOW, C.DATE_CONFIRM)
rect(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(3.2), fill=MIST)
rect(s, Inches(0.5), Inches(1.5), Inches(6.0), Pt(5), fill=TEAL)
textbox(s, Inches(0.75), Inches(1.75), Inches(5.5), Inches(2.7),
        [(C.TIME_OPTIONS[0][0], 16, NAVY, True, 8),
         (C.TIME_OPTIONS[0][1], 13, TEAL, True, 8),
         (C.TIME_OPTIONS[0][2], 13, DARK, False, 10),
         ("穿山港区 + 航运展厅 + 海天智造", 12, GREY, False)])

rect(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(3.2), fill=NAVY)
textbox(s, Inches(7.05), Inches(1.75), Inches(5.5), Inches(2.7),
        [(C.TIME_OPTIONS[1][0], 16, AMBER, True, 8),
         (C.TIME_OPTIONS[1][1], 13, SEA, True, 8),
         (C.TIME_OPTIONS[1][2], 13, WHITE, False, 10),
         ("加极氪/硬科技深访 + 孵化落地 + 投洽", 12, RGBColor(0xB8, 0xCF, 0xDC), False)])

textbox(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.6),
        [("建议窗口", 13, TEAL, True, 6)] +
        [(f"· {m}", 12, DARK, False, 5) for m in C.PREFERRED_MONTHS])
footer(s, 11)

# ================================================================ 12. 第一日行程
s = add_slide()
header(s, "05 行程 · 第一日", "已嵌入港口精选点位与互动链接一、二", "一日行至此结束")
for i, (t, act, note) in enumerate(C.DAY1_ITINERARY):
    y = Inches(1.35) + i * Inches(0.65)
    bg = MIST if i % 2 == 0 else WHITE
    rect(s, Inches(0.45), y, Inches(12.4), Inches(0.58), fill=bg)
    textbox(s, Inches(0.55), y + Inches(0.08), Inches(2.3), Inches(0.42),
            [(t, 11, TEAL, True)], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(2.9), y + Inches(0.08), Inches(4.2), Inches(0.42),
            [(act, 11.5, NAVY, True)], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(7.2), y + Inches(0.08), Inches(5.4), Inches(0.42),
            [(note, 11, GREY, False)], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 12)

# ================================================================ 13. 第二日行程
s = add_slide()
header(s, "05 行程 · 第二日", "孵化落地 · 新能源智造 · 硬科技 · 投洽", "推荐两日深度行")
for i, (t, act, note) in enumerate(C.DAY2_ITINERARY):
    y = Inches(1.35) + i * Inches(0.72)
    rect(s, Inches(0.5), y, Inches(12.3), Inches(0.65), fill=MIST if i % 2 == 0 else LIGHT)
    rect(s, Inches(0.5), y, Pt(5), Inches(0.65), fill=AMBER if i == 4 else TEAL)
    textbox(s, Inches(0.75), y + Inches(0.1), Inches(2.4), Inches(0.45),
            [(t, 11, TEAL, True)], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(3.3), y + Inches(0.1), Inches(4.5), Inches(0.45),
            [(act, 12, NAVY, True)], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(7.9), y + Inches(0.1), Inches(4.7), Inches(0.45),
            [(note, 11, GREY, False)], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 13)

# ================================================================ 14. 成果与长效
s = add_slide()
header(s, "06 预期成果", "四类线索清单 · 可跟进可转化", "")
textbox(s, Inches(0.5), Inches(1.4), Inches(6.2), Inches(5.2),
        [("本次预期成果", 14, TEAL, True, 8)] +
        [(f"· {o}", 12.5, DARK, False, 10) for o in C.OUTCOMES])
rect(s, Inches(6.9), Inches(1.4), Inches(5.9), Inches(5.2), fill=DEEP)
textbox(s, Inches(7.2), Inches(1.65), Inches(5.4), Inches(0.4),
        [("长效合作设想", 14, AMBER, True)])
lines = []
for name, desc in C.LONG_TERM:
    lines.append((name, 13, SEA, True, 4))
    lines.append((desc, 12, WHITE, False, 10))
textbox(s, Inches(7.2), Inches(2.2), Inches(5.4), Inches(4.1), lines)
footer(s, 14)

# ================================================================ 15. 恳请支持
s = add_slide()
header(s, "06 恳请支持", "请贵局协助协调港口、企业与互动环节", "")
for i, (title, desc) in enumerate(C.SUPPORT_REQUESTS):
    y = Inches(1.32) + i * Inches(0.75)
    rect(s, Inches(0.5), y, Inches(12.3), Inches(0.68), fill=MIST if i % 2 == 0 else LIGHT)
    textbox(s, Inches(0.7), y + Inches(0.12), Inches(2.8), Inches(0.45),
            [(f"{i+1}. {title}", 13, NAVY, True)], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(3.6), y + Inches(0.12), Inches(9.0), Inches(0.45),
            [(desc, 12, DARK, False)], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 15)

# ================================================================ 16. 下一步
s = add_slide()
header(s, "下一步", "从策划案到成行：先锁定点位再确认人员", "")
for i, (step, work) in enumerate(C.NEXT_STEPS):
    y = Inches(1.5) + i * Inches(0.95)
    rect(s, Inches(0.5), y, Inches(12.3), Inches(0.85), fill=MIST if i % 2 == 0 else LIGHT)
    textbox(s, Inches(0.7), y + Inches(0.2), Inches(3.5), Inches(0.45),
            [(step, 14, TEAL, True)], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(4.4), y + Inches(0.2), Inches(8.1), Inches(0.45),
            [(work, 13, DARK, False)], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 16)

# ================================================================ 17. 封底
s = add_slide()
rect(s, 0, 0, SW, SH, fill=DEEP)
rect(s, 0, 0, Inches(0.18), SH, fill=TEAL)
textbox(s, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.5),
        [("期待与贵局携手", 14, SEA, True)])
textbox(s, Inches(0.8), Inches(2.4), Inches(11.5), Inches(1.5),
        [("精选港口点位 · 链接产业与资本", 26, WHITE, True, 10),
         ("推动沪甬可跟进、可转化的务实合作", 26, WHITE, True)])
textbox(s, Inches(0.8), Inches(4.5), Inches(11.5), Inches(1.0),
        [("港口枢纽 → 港航服务 → 临港智造 → 科创落地 → 资本投洽", 14,
          RGBColor(0xB8, 0xCF, 0xDC), False)])
textbox(s, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.0),
        [("复旦大学住房政策研究中心", 14, WHITE, True, 6),
         ("上海市杨浦区科技企业联合会", 14, WHITE, True, 6),
         ("提交：" + C.RECIPIENT + "  ·  " + C.DOC_DATE + "  ·  " + C.VERSION, 12, AMBER, False)])

OUT_DIR = os.path.join(os.path.dirname(__file__), "output")
os.makedirs(OUT_DIR, exist_ok=True)
out = os.path.join(OUT_DIR, "赴宁波港及宁波经济技术开发区考察交流策划案.pptx")
prs.save(out)
print("saved:", out)
